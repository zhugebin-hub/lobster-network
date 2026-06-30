#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
人教版高中物理必修 2《动能 动能定理》PPT 课件生成器
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def create_kinetic_energy_ppt():
    # 创建演示文稿
    prs = Presentation()
    
    # 设置宽屏 16:9
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # 定义颜色
    BLUE = RGBColor(52, 101, 164)
    LIGHT_BLUE = RGBColor(230, 240, 250)
    DARK_BLUE = RGBColor(30, 60, 100)
    ORANGE = RGBColor(245, 124, 0)
    GREEN = RGBColor(67, 160, 71)
    GRAY = RGBColor(100, 100, 100)
    
    def add_title_slide(prs, title, subtitle):
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白版式
        
        # 背景色块
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.5)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = BLUE
        shape.line.fill.background()
        
        # 标题
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        
        # 副标题
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(12), Inches(1))
        tf = subtitle_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.font.color.rgb = GRAY
        p.alignment = PP_ALIGN.CENTER
        
        return slide
    
    def add_content_slide(prs, title, content_items, notes=""):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # 标题栏
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = BLUE
        
        # 下划线
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.1), Inches(2), Inches(0.05)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = ORANGE
        line.line.fill.background()
        
        # 内容
        y_pos = 1.5
        for item in content_items:
            if isinstance(item, dict):
                # 带标题的内容块
                box = slide.shapes.add_textbox(Inches(0.7), Inches(y_pos), Inches(12), Inches(0.8))
                tf = box.text_frame
                p = tf.paragraphs[0]
                p.text = item.get('title', '')
                p.font.size = Pt(24)
                p.font.bold = True
                p.font.color.rgb = DARK_BLUE
                y_pos += 0.7
                
                for subitem in item.get('items', []):
                    box = slide.shapes.add_textbox(Inches(1), Inches(y_pos), Inches(11.5), Inches(0.5))
                    tf = box.text_frame
                    p = tf.paragraphs[0]
                    p.text = subitem
                    p.font.size = Pt(20)
                    p.font.color.rgb = GRAY
                    y_pos += 0.5
                y_pos += 0.2
            else:
                box = slide.shapes.add_textbox(Inches(0.7), Inches(y_pos), Inches(12), Inches(0.6))
                tf = box.text_frame
                p = tf.paragraphs[0]
                p.text = item
                p.font.size = Pt(22)
                p.font.color.rgb = GRAY
                y_pos += 0.6
        
        return slide
    
    def add_formula_slide(prs, title, formula, explanation, derivation_steps=None):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # 标题
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = BLUE
        
        # 公式框（突出显示）
        formula_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(1.3), Inches(11.3), Inches(2)
        )
        formula_box.fill.solid()
        formula_box.fill.fore_color.rgb = LIGHT_BLUE
        formula_box.line.color.rgb = BLUE
        formula_box.line.width = Pt(3)
        
        # 公式内容
        tf = formula_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = formula
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = DARK_BLUE
        p.alignment = PP_ALIGN.CENTER
        
        # 说明
        y_pos = 3.6
        for line in explanation.split('\n'):
            box = slide.shapes.add_textbox(Inches(0.7), Inches(y_pos), Inches(12), Inches(0.5))
            tf = box.text_frame
            p = tf.paragraphs[0]
            p.text = line
            p.font.size = Pt(20)
            p.font.color.rgb = GRAY
            y_pos += 0.5
        
        # 推导步骤（如果有）
        if derivation_steps:
            y_pos = 5.0
            for i, step in enumerate(derivation_steps, 1):
                box = slide.shapes.add_textbox(Inches(1), Inches(y_pos), Inches(11.5), Inches(0.4))
                tf = box.text_frame
                p = tf.paragraphs[0]
                p.text = f"{i}. {step}"
                p.font.size = Pt(18)
                p.font.color.rgb = GRAY
                y_pos += 0.4
        
        return slide
    
    def add_diagram_slide(prs, title, diagram_description):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # 标题
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = BLUE
        
        # 示意图区域
        diagram_box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(1), Inches(1.3), Inches(11.3), Inches(4.5)
        )
        diagram_box.fill.solid()
        diagram_box.fill.fore_color.rgb = LIGHT_BLUE
        diagram_box.line.color.rgb = BLUE
        diagram_box.line.width = Pt(2)
        
        # 示意图说明文字
        tf = diagram_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = diagram_description
        p.font.size = Pt(18)
        p.font.color.rgb = GRAY
        p.alignment = PP_ALIGN.CENTER
        
        return slide
    
    def add_animation_sequence_slide(prs, title, steps):
        """创建动画演示序列幻灯片"""
        slides = []
        for i, step in enumerate(steps):
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            
            # 标题（带步骤指示）
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
            tf = title_box.text_frame
            p = tf.paragraphs[0]
            p.text = f"{title} - 步骤 {i+1}/{len(steps)}"
            p.font.size = Pt(32)
            p.font.bold = True
            p.font.color.rgb = BLUE
            
            # 步骤内容
            y_pos = 1.3
            for line in step.split('\n'):
                box = slide.shapes.add_textbox(Inches(0.7), Inches(y_pos), Inches(12), Inches(0.6))
                tf = box.text_frame
                p = tf.paragraphs[0]
                p.text = line
                p.font.size = Pt(22)
                p.font.color.rgb = GRAY
                y_pos += 0.6
            
            slides.append(slide)
        
        return slides
    
    # ========== 开始创建幻灯片 ==========
    
    # 第 1 页：封面
    add_title_slide(prs, "动能 动能定理", "人教版高中物理必修 2 · 第七章 机械能守恒定律")
    
    # 第 2 页：学习目标
    add_content_slide(prs, "🎯 学习目标", [
        {"title": "知识与技能", "items": [
            "理解动能的概念，掌握动能的表达式",
            "理解动能定理的内容和物理意义",
            "能运用动能定理解决简单的力学问题"
        ]},
        {"title": "过程与方法", "items": [
            "通过理论推导，探究动能定理的建立过程",
            "体会演绎推理在物理研究中的应用"
        ]},
        {"title": "情感态度与价值观", "items": [
            "感受物理规律的简洁美",
            "培养严谨的科学态度"
        ]}
    ])
    
    # 第 3 页：情境导入
    add_content_slide(prs, "💡 情境导入", [
        "思考以下问题：",
        "• 为什么高速行驶的汽车刹车距离更长？",
        "• 为什么子弹虽然质量小，但杀伤力很大？",
        "• 物体的动能与哪些因素有关？",
        "",
        "【生活实例】",
        "• 流星撞击地球释放巨大能量",
        "• 风力发电机利用空气动能发电",
        "• 打桩机利用重锤的动能做功"
    ])
    
    # 第 4 页：动能的概念
    add_content_slide(prs, "📖 一、动能的概念", [
        {"title": "定义", "items": [
            "物体由于运动而具有的能量叫做动能"
        ]},
        {"title": "表达式", "items": [
            "Ek = ½mv²",
            "其中：m — 物体的质量（kg）",
            "      v — 物体的瞬时速度（m/s）",
            "      Ek — 动能（J）"
        ]},
        {"title": "单位", "items": [
            "国际单位：焦耳（J）",
            "1 J = 1 kg·m²/s² = 1 N·m"
        ]},
        {"title": "特点", "items": [
            "• 动能是标量，只有大小，没有方向",
            "• 动能总是正值（v²≥0）",
            "• 动能具有相对性，与参考系的选择有关"
        ]}
    ])
    
    # 第 5 页：动能表达式推导
    add_formula_slide(prs, "📐 动能表达式的推导", 
        "Ek = ½mv²",
        "推导思路：从功与能的关系出发",
        [
            "设质量为 m 的物体，初速度为 v₀，在恒力 F 作用下",
            "经过位移 l，末速度为 v",
            "由牛顿第二定律：F = ma",
            "由运动学公式：v² - v₀² = 2al",
            "力做的功：W = Fl = ma·l = m·(v²-v₀²)/2",
            "即：W = ½mv² - ½mv₀²"
        ]
    )
    
    # 第 6-9 页：动能定理推导动画序列
    anim_steps = [
        "【步骤 1】物理模型\n\n• 光滑水平面上，质量为 m 的物体\n• 初速度 v₀，受恒力 F 作用\n• 经过位移 l，末速度为 v\n\n→ 分析力做功与速度变化的关系",
        
        "【步骤 2】牛顿第二定律\n\n• 物体受力：F（合力）\n• 加速度：a = F/m\n• 物体做匀加速直线运动\n\n→ 建立力与加速度的关系",
        
        "【步骤 3】运动学公式\n\n• 匀变速直线运动：v² - v₀² = 2al\n• 变形得：l = (v² - v₀²)/(2a)\n\n→ 建立位移与速度的关系",
        
        "【步骤 4】功的计算\n\n• 合力做功：W = F·l\n• 代入 l：W = F·(v² - v₀²)/(2a)\n• 代入 a = F/m：W = ½mv² - ½mv₀²\n\n→ 得到功与动能变化的关系"
    ]
    
    for i, step in enumerate(anim_steps):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # 标题
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"🎬 动能定理推导演示 - 步骤 {i+1}/4"
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = BLUE
        
        # 进度条
        progress = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.0), Inches(12.3), Inches(0.1)
        )
        progress.fill.solid()
        progress.fill.fore_color.rgb = LIGHT_BLUE
        progress.line.fill.background()
        
        progress_fill = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.0), Inches(12.3 * (i+1) / 4), Inches(0.1)
        )
        progress_fill.fill.solid()
        progress_fill.fill.fore_color.rgb = ORANGE
        progress_fill.line.fill.background()
        
        # 内容
        y_pos = 1.3
        for line in step.split('\n'):
            if line.strip():
                box = slide.shapes.add_textbox(Inches(0.7), Inches(y_pos), Inches(12), Inches(0.5))
                tf = box.text_frame
                p = tf.paragraphs[0]
                p.text = line
                p.font.size = Pt(20)
                p.font.color.rgb = GRAY
                y_pos += 0.5
    
    # 第 10 页：动能定理
    add_formula_slide(prs, "📌 二、动能定理", 
        "W = Ek₂ - Ek₁ = ΔEk",
        "内容：合力对物体所做的功，等于物体动能的变化量\n\n物理意义：\n• 揭示了功与能的转化关系\n• 功是能量转化的量度\n• 动能的变化由合力做功决定",
        [
            "W > 0：合力做正功，动能增加",
            "W < 0：合力做负功，动能减少",
            "W = 0：合力不做功，动能不变"
        ]
    )
    
    # 第 11 页：动能定理的理解
    add_content_slide(prs, "💭 动能定理的理解", [
        {"title": "适用范围", "items": [
            "• 既适用于恒力做功，也适用于变力做功",
            "• 既适用于直线运动，也适用于曲线运动",
            "• 既适用于单个物体，也适用于物体系"
        ]},
        {"title": "W 的含义", "items": [
            "• W 是合力做的功（或各力做功的代数和）",
            "• W = W₁ + W₂ + W₃ + ...",
            "• 包括重力、弹力、摩擦力等所有力"
        ]},
        {"title": "解题优势", "items": [
            "• 不涉及加速度和时间，简化计算",
            "• 特别适用于变力做功和曲线运动问题",
            "• 是解决力学问题的重要工具"
        ]}
    ])
    
    # 第 12 页：应用示例 1
    add_content_slide(prs, "✏️ 应用示例 1：刹车距离", [
        {"title": "题目", "items": [
            "一辆质量 m = 1500 kg 的汽车，以 v = 20 m/s 的速度行驶",
            "刹车时受到的阻力 f = 6000 N",
            "求：汽车的刹车距离"
        ]},
        {"title": "解析", "items": [
            "由动能定理：W = ΔEk",
            "-f·s = 0 - ½mv²",
            "s = mv²/(2f) = 1500×20²/(2×6000) = 50 m",
            "",
            "【答案】刹车距离为 50 米"
        ]},
        {"title": "思考", "items": [
            "• 速度加倍，刹车距离变为多少？（4 倍！）",
            "• 这就是为什么不能超速行驶"
        ]}
    ])
    
    # 第 13 页：应用示例 2
    add_content_slide(prs, "✏️ 应用示例 2：自由落体", [
        {"title": "题目", "items": [
            "质量为 m 的物体从高度 h 处自由下落",
            "求：落地时的速度",
            "（不计空气阻力）"
        ]},
        {"title": "解析", "items": [
            "由动能定理：W = ΔEk",
            "mgh = ½mv² - 0",
            "v = √(2gh)",
            "",
            "【答案】落地速度 v = √(2gh)"
        ]},
        {"title": "说明", "items": [
            "• 与运动学公式结果一致",
            "• 动能定理更简洁，不涉及时间"
        ]}
    ])
    
    # 第 14 页：课堂练习
    add_content_slide(prs, "📝 课堂练习", [
        {"title": "基础题", "items": [
            "1. 质量为 2 kg 的物体，速度从 3 m/s 增加到 5 m/s",
            "   求动能的变化量",
            "   【答案】ΔEk = ½×2×(5²-3²) = 16 J"
        ]},
        {"title": "提升题", "items": [
            "2. 物体从斜面顶端滑下，高度 h = 5 m，斜面长 l = 10 m",
            "   摩擦因数μ = 0.2，求到达底端的速度",
            "   【提示】mgh - μmgcosθ·l = ½mv²"
        ]},
        {"title": "拓展题", "items": [
            "3. 思考：为什么过山车要从高处开始？",
            "   用动能定理解释"
        ]}
    ])
    
    # 第 15 页：易错点警示
    add_content_slide(prs, "⚠️ 易错点警示", [
        {"title": "常见错误", "items": [
            "❌ 忘记动能是标量，误认为有方向",
            "❌ 混淆动能和动量（mv）",
            "❌ 计算合力功时漏掉某个力",
            "❌ 动能定理中 W 是合力功，不是某个力的功"
        ]},
        {"title": "注意事项", "items": [
            "✓ 明确研究对象和过程",
            "✓ 正确分析受力，计算各力做功",
            "✓ 注意初末状态的动能",
            "✓ 单位统一用国际单位制"
        ]}
    ])
    
    # 第 16 页：知识总结
    add_content_slide(prs, "📊 知识总结", [
        {"title": "核心公式", "items": [
            "动能：Ek = ½mv²",
            "动能定理：W = ΔEk = Ek₂ - Ek₁"
        ]},
        {"title": "物理思想", "items": [
            "• 能量观：功是能量转化的量度",
            "• 状态观：动能是状态量，功是过程量",
            "• 转化观：不同形式能量可以相互转化"
        ]},
        {"title": "解题步骤", "items": [
            "1. 确定研究对象和过程",
            "2. 分析受力，计算合力做功",
            "3. 确定初末动能",
            "4. 列动能定理方程求解"
        ]}
    ])
    
    # 第 17 页：课后作业
    add_content_slide(prs, "📚 课后作业", [
        "【必做题】教材 P73 练习 1、2、3",
        "",
        "【选做题】",
        "• 查阅资料：动能定理在生活中的应用",
        "• 思考题：如果考虑空气阻力，自由落体的",
        "  落地速度如何计算？",
        "",
        "【预习】",
        "• 预习下一节：机械能守恒定律",
        "• 思考：什么情况下机械能守恒？"
    ])
    
    # 第 18 页：结束页
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 背景
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, Inches(2), prs.slide_width, Inches(5.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = BLUE
    shape.line.fill.background()
    
    # 文字
    box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.3), Inches(2))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = "谢谢观看！"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    p2 = tf.add_paragraph()
    p2.text = "勤学善思 · 格物致知"
    p2.font.size = Pt(24)
    p2.font.color.rgb = RGBColor(255, 255, 255)
    p2.alignment = PP_ALIGN.CENTER
    
    # 保存文件
    output_path = "/home/admin/.openclaw/workspace/动能_动能定理.pptx"
    prs.save(output_path)
    
    return output_path

if __name__ == "__main__":
    path = create_kinetic_energy_ppt()
    print(f"PPT 已生成：{path}")
