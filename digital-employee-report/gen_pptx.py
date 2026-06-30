#!/usr/bin/env python3
"""数字员工落地汇报 PPT 生成脚本"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ─── Color palette ───
BG_DARK = RGBColor(0x0A, 0x16, 0x28)      # 深蓝黑背景
BG_CARD = RGBColor(0x0F, 0x28, 0x47)       # 卡片背景
ACCENT = RGBColor(0x3B, 0x82, 0xF6)        # 主色调蓝色
ACCENT_LIGHT = RGBColor(0x60, 0xA5, 0xFA)  # 亮蓝
ACCENT_PURPLE = RGBColor(0xA7, 0x8B, 0xFA) # 紫色渐变
MUTED = RGBColor(0x7A, 0x8B, 0xA3)         # 灰色文字
WHITE = RGBColor(0xE8, 0xEC, 0xF1)         # 主文字
GREEN = RGBColor(0x10, 0xB9, 0x81)         # 绿色指标
YELLOW = RGBColor(0xF5, 0x9E, 0x0B)        # 黄色
ORANGE = RGBColor(0xF5, 0x9E, 0x0B)
PURPLE = RGBColor(0x8B, 0x5C, 0xF6)
WHITE_PURE = RGBColor(0xFF, 0xFF, 0xFF)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_WIDTH
prs.slide_height = SLIDE_HEIGHT

# ─── Helpers ───

def add_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape_bg(slide, left, top, width, height, color, transparency=0):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_rounded_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_circle(slide, left, top, size, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name='微软雅黑'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_multi_text(slide, left, top, width, height, lines, font_name='微软雅黑'):
    """lines: [(text, font_size, color, bold, alignment), ...] alignment defaults to LEFT"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(lines):
        text = item[0]
        font_size = item[1] if len(item) > 1 else 18
        color = item[2] if len(item) > 2 else WHITE
        bold = item[3] if len(item) > 3 else False
        alignment = item[4] if len(item) > 4 else PP_ALIGN.LEFT
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font_name
        p.alignment = alignment
        p.space_after = Pt(8)
    return txBox

# ════════════════════════════════════════════════════════════
# SLIDE 1: COVER
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, BG_DARK)
add_shape_bg(slide, Inches(0), Inches(0), SLIDE_WIDTH, SLIDE_HEIGHT, BG_DARK)

# Decorative circles
add_circle(slide, Inches(9), Inches(-0.5), Inches(4), RGBColor(0x0F, 0x28, 0x47))
add_circle(slide, Inches(-1), Inches(5.5), Inches(3), RGBColor(0x0F, 0x28, 0x47))

# Top accent line
add_shape_bg(slide, Inches(3), Inches(1.2), Inches(7.333), Inches(0.03), ACCENT)

# Main title
add_textbox(slide, Inches(1.5), Inches(2.0), Inches(10.333), Inches(1.5),
            "数字员工落地实践", font_size=52, color=ACCENT_LIGHT, bold=True,
            alignment=PP_ALIGN.CENTER)

# Subtitle
add_textbox(slide, Inches(2), Inches(3.5), Inches(9.333), Inches(1.2),
            "从规划到落地 · 构建智能化组织，释放AI生产力", font_size=28,
            color=MUTED, alignment=PP_ALIGN.CENTER)

# Bottom info
add_shape_bg(slide, Inches(3), Inches(5.8), Inches(7.333), Inches(0.03), ACCENT)
add_textbox(slide, Inches(2), Inches(6.1), Inches(9.333), Inches(0.8),
            "2026年6月", font_size=22, color=MUTED, alignment=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════
# SLIDE 2: AGENDA
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)

add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11.5), Inches(1),
            "汇报目录", font_size=44, color=WHITE, bold=True)

# 6 items in 2x3 grid
agenda_items = [
    ("01", "背景与趋势", "数字化转型浪潮与AI技术突破"),
    ("02", "数字员工定义", "能力边界与技术架构"),
    ("03", "落地方案设计", "整体架构与实施路径"),
    ("04", "典型应用场景", "客服、运营、数据分析等"),
    ("05", "成本与收益", "ROI分析与投资回报"),
    ("06", "下一步计划", "实施路线图与里程碑"),
]

for i, (num, title, desc) in enumerate(agenda_items):
    col = i % 3
    row = i // 3
    left = Inches(0.8 + col * 4.1)
    top = Inches(1.8 + row * 2.5)

    card = add_rounded_rect(slide, left, top, Inches(3.7), Inches(2.1), BG_CARD)
    card.line.color.rgb = RGBColor(0x1A, 0x3A, 0x6A)
    card.line.width = Pt(1)

    add_textbox(slide, left + Inches(0.3), top + Inches(0.3), Inches(1.2), Inches(0.8),
                num, font_size=40, color=ACCENT, bold=True)
    add_textbox(slide, left + Inches(1.6), top + Inches(0.35), Inches(1.8), Inches(0.6),
                title, font_size=28, color=WHITE, bold=True)
    add_textbox(slide, left + Inches(0.3), top + Inches(1.2), Inches(3.1), Inches(0.8),
                desc, font_size=20, color=MUTED)

# ════════════════════════════════════════════════════════════
# SLIDE 3: SECTION DIVIDER - Chapter 01
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)
add_shape_bg(slide, Inches(0), Inches(0), SLIDE_WIDTH, SLIDE_HEIGHT,
             RGBColor(0x0F, 0x28, 0x47))

add_textbox(slide, Inches(2), Inches(2.2), Inches(9.333), Inches(0.8),
            "CHAPTER 01", font_size=24, color=ACCENT, bold=True,
            alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1.5), Inches(3.2), Inches(10.333), Inches(1.5),
            "背景与趋势", font_size=56, color=WHITE, bold=True,
            alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(2), Inches(4.8), Inches(9.333), Inches(0.8),
            "为什么现在是数字员工的最佳时机", font_size=26, color=MUTED,
            alignment=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════
# SLIDE 4: Background - Industry Trends
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)

add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11.5), Inches(0.8),
            "01 行业背景", font_size=40, color=WHITE, bold=True)

# 3 stat cards
stats = [
    ("78%", "企业计划部署AI员工", "Gartner 2026"),
    ("¥2.3万亿", "中国AI市场规模", "IDC预测"),
    ("40%+", "人力成本可降低", "麦肯锡报告"),
]

for i, (val, label, src) in enumerate(stats):
    left = Inches(0.8 + i * 4.1)
    card = add_rounded_rect(slide, left, Inches(1.3), Inches(3.7), Inches(2.8), BG_CARD)
    card.line.color.rgb = RGBColor(0x1A, 0x3A, 0x6A)
    card.line.width = Pt(1)

    add_textbox(slide, left + Inches(0.3), Inches(1.6), Inches(3.1), Inches(0.8),
                val, font_size=42, color=ACCENT_LIGHT, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, left + Inches(0.3), Inches(2.5), Inches(3.1), Inches(0.8),
                label, font_size=24, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, left + Inches(0.3), Inches(3.3), Inches(3.1), Inches(0.6),
                src, font_size=18, color=MUTED, alignment=PP_ALIGN.CENTER)

# Bottom quote box
quote_box = add_rounded_rect(slide, Inches(0.8), Inches(4.5), Inches(11.7), Inches(2.5),
                             RGBColor(0x0F, 0x28, 0x47))
quote_box.line.color.rgb = ACCENT
quote_box.line.width = Pt(2)

add_multi_text(slide, Inches(1.2), Inches(4.7), Inches(10.9), Inches(2.1), [
    ("💡 核心判断", 28, ACCENT, True, PP_ALIGN.LEFT),
    ("大语言模型能力突破 + 自动化技术成熟 + 企业数字化基础完备", 26, WHITE, False, PP_ALIGN.LEFT),
    ("= 数字员工规模化落地窗口已开启", 26, ACCENT_LIGHT, True, PP_ALIGN.LEFT),
])

# ════════════════════════════════════════════════════════════
# SLIDE 5: SECTION DIVIDER - Chapter 02
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, RGBColor(0x0F, 0x28, 0x47))

add_textbox(slide, Inches(2), Inches(2.2), Inches(9.333), Inches(0.8),
            "CHAPTER 02", font_size=24, color=ACCENT, bold=True,
            alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1.5), Inches(3.2), Inches(10.333), Inches(1.5),
            "数字员工定义", font_size=56, color=WHITE, bold=True,
            alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(2), Inches(4.8), Inches(9.333), Inches(0.8),
            "能力边界与技术架构", font_size=26, color=MUTED,
            alignment=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════
# SLIDE 6: What is Digital Employee
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)

add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11.5), Inches(0.8),
            "02 什么是数字员工？", font_size=40, color=WHITE, bold=True)

# Definition card
def_card = add_rounded_rect(slide, Inches(0.8), Inches(1.2), Inches(7.2), Inches(1.8), BG_CARD)
def_card.line.color.rgb = RGBColor(0x1A, 0x3A, 0x6A)
def_card.line.width = Pt(1)

add_multi_text(slide, Inches(1.2), Inches(1.3), Inches(6.5), Inches(1.6), [
    ("核心定义", 22, ACCENT, True),
    ("基于AI技术的软件机器人，能够模拟人类员工的认知能力，", 22, WHITE),
    ("独立完成特定业务流程中的工作任务。", 22, WHITE),
])

# 4 capabilities
caps = [
    ("理解能力", "自然语言处理、意图识别、上下文理解"),
    ("决策能力", "规则引擎、智能判断、异常处理"),
    ("执行能力", "API调用、系统操作、多平台协作"),
    ("学习能力", "持续优化、知识更新、自我迭代"),
]
for i, (title, desc) in enumerate(caps):
    top = Inches(3.4 + i * 0.85)
    add_circle(slide, Inches(1.0), top + Inches(0.15), Inches(0.25), ACCENT)
    add_textbox(slide, Inches(1.4), top, Inches(1.5), Inches(0.5),
                title, font_size=22, color=WHITE, bold=True)
    add_textbox(slide, Inches(3.0), top, Inches(5.0), Inches(0.5),
                desc, font_size=20, color=MUTED)

# Right: vs RPA comparison
vs_card = add_rounded_rect(slide, Inches(8.5), Inches(1.2), Inches(4.2), Inches(5.5),
                           RGBColor(0x0F, 0x28, 0x47))
vs_card.line.color.rgb = RGBColor(0x1A, 0x3A, 0x6A)
vs_card.line.width = Pt(1)

add_textbox(slide, Inches(8.7), Inches(1.3), Inches(3.8), Inches(0.6),
            "vs 传统RPA", font_size=26, color=WHITE, bold=True,
            alignment=PP_ALIGN.CENTER)

vs_items = [
    ("规则驱动", "语义驱动"),
    ("固定流程", "灵活适应"),
    ("结构化数据", "非结构化数据"),
    ("需要人工配置", "自主学习"),
]
for i, (old, new) in enumerate(vs_items):
    top = Inches(2.1 + i * 1.0)
    add_textbox(slide, Inches(8.7), top, Inches(1.8), Inches(0.5),
                old, font_size=20, color=MUTED, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(10.8), top, Inches(1.8), Inches(0.5),
                new, font_size=20, color=ACCENT_LIGHT, bold=True,
                alignment=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════
# SLIDE 7: Architecture
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)

add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11.5), Inches(0.8),
            "02 技术架构", font_size=40, color=WHITE, bold=True)

layers = [
    ("交互层", "自然语言对话  |  多模态输入  |  语音/文字/图像", GREEN),
    ("大脑层", "大语言模型 (LLM)  |  知识图谱  |  推理与规划引擎", ACCENT),
    ("能力层", "API调用  |  代码执行  |  文档处理  |  数据分析", PURPLE),
    ("执行层", "RPA自动化  |  业务流程编排  |  多系统对接", YELLOW),
]
for i, (layer, items, color) in enumerate(layers):
    top = Inches(1.3 + i * 1.5)
    # Layer label
    lbl = add_rounded_rect(slide, Inches(0.8), top, Inches(1.5), Inches(1.1),
                           RGBColor(0x0F, 0x28, 0x47))
    lbl.line.color.rgb = color
    lbl.line.width = Pt(2)
    add_textbox(slide, Inches(0.8), top + Inches(0.3), Inches(1.5), Inches(0.5),
                layer, font_size=22, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    # Items
    add_rounded_rect(slide, Inches(2.5), top, Inches(10), Inches(1.1), BG_CARD)
    add_textbox(slide, Inches(2.7), top + Inches(0.25), Inches(9.5), Inches(0.6),
                items, font_size=22, color=WHITE)

# ════════════════════════════════════════════════════════════
# SLIDE 8: SECTION DIVIDER - Chapter 03
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, RGBColor(0x0F, 0x28, 0x47))

add_textbox(slide, Inches(2), Inches(2.2), Inches(9.333), Inches(0.8),
            "CHAPTER 03", font_size=24, color=ACCENT, bold=True,
            alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1.5), Inches(3.2), Inches(10.333), Inches(1.5),
            "落地方案设计", font_size=56, color=WHITE, bold=True,
            alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(2), Inches(4.8), Inches(9.333), Inches(0.8),
            "整体架构与实施路径", font_size=26, color=MUTED,
            alignment=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════
# SLIDE 9: Implementation Roadmap
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)

add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11.5), Inches(0.8),
            "03 三阶段实施路径", font_size=40, color=WHITE, bold=True)

phases = [
    ("第一阶段", "1-3个月", "试点验证", GREEN,
     ["选取1-2个高频场景", "部署基础版数字员工", "跑通业务流程闭环", "评估效果并优化"]),
    ("第二阶段", "3-6个月", "规模推广", ACCENT,
     ["扩展至5-8个场景", "建立数字员工管理平台", "培训业务人员协同", "形成标准化SOP"]),
    ("第三阶段", "6-12个月", "全面融合", PURPLE,
     ["覆盖核心业务线", "数字员工自主运营", "人机协同常态化", "持续迭代优化"]),
]

for i, (phase, time, tag, color, items) in enumerate(phases):
    left = Inches(0.8 + i * 4.1)
    card = add_rounded_rect(slide, left, Inches(1.3), Inches(3.7), Inches(5.5), BG_CARD)
    card.line.color.rgb = color
    card.line.width = Pt(2)

    # Phase header
    add_textbox(slide, left + Inches(0.3), Inches(1.5), Inches(1.8), Inches(0.4),
                phase, font_size=20, color=color, bold=True)
    add_textbox(slide, left + Inches(2.2), Inches(1.5), Inches(1.2), Inches(0.4),
                time, font_size=18, color=MUTED)
    add_textbox(slide, left + Inches(0.3), Inches(2.0), Inches(3.1), Inches(0.6),
                tag, font_size=30, color=WHITE, bold=True)

    # Items
    for j, item in enumerate(items):
        top = Inches(2.9 + j * 0.9)
        add_circle(slide, left + Inches(0.3), top + Inches(0.15), Inches(0.2), color)
        add_textbox(slide, left + Inches(0.7), top, Inches(2.7), Inches(0.5),
                    item, font_size=20, color=MUTED)

# ════════════════════════════════════════════════════════════
# SLIDE 10: Key Success Factors
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)

add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11.5), Inches(0.8),
            "03 成功关键要素", font_size=40, color=WHITE, bold=True)

factors = [
    ("🎯", "明确业务目标", "从痛点出发，而非技术驱动。先选高频、低风险的场景验证价值。"),
    ("🔗", "高层支持与跨部门协作", "一把手工程，IT与业务部门紧密配合，确保数据与流程打通。"),
    ("📊", "可量化指标体系", "建立ROI、效率提升、错误率下降等量化指标，持续跟踪。"),
    ("🔒", "安全与合规保障", "数据隐私保护、权限管控、操作审计，确保合规运行。"),
]
for i, (icon, title, desc) in enumerate(factors):
    col = i % 2
    row = i // 2
    left = Inches(0.8 + col * 6)
    top = Inches(1.3 + row * 2.9)

    card = add_rounded_rect(slide, left, top, Inches(5.6), Inches(2.5), BG_CARD)
    card.line.color.rgb = RGBColor(0x1A, 0x3A, 0x6A)
    card.line.width = Pt(1)

    add_textbox(slide, left + Inches(0.3), top + Inches(0.2), Inches(5), Inches(0.5),
                f"{icon} {title}", font_size=26, color=WHITE, bold=True)
    add_textbox(slide, left + Inches(0.3), top + Inches(0.8), Inches(5), Inches(1.5),
                desc, font_size=20, color=MUTED)

# ════════════════════════════════════════════════════════════
# SLIDE 11: SECTION DIVIDER - Chapter 04
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, RGBColor(0x0F, 0x28, 0x47))

add_textbox(slide, Inches(2), Inches(2.2), Inches(9.333), Inches(0.8),
            "CHAPTER 04", font_size=24, color=ACCENT, bold=True,
            alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1.5), Inches(3.2), Inches(10.333), Inches(1.5),
            "典型应用场景", font_size=56, color=WHITE, bold=True,
            alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(2), Inches(4.8), Inches(9.333), Inches(0.8),
            "客服、运营、数据分析等", font_size=26, color=MUTED,
            alignment=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════
# SLIDE 12: Use Cases
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)

add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11.5), Inches(0.8),
            "04 核心应用场景", font_size=40, color=WHITE, bold=True)

use_cases = [
    ("💬", "智能客服", GREEN, "响应时间↓80%",
     ["7×24在线应答", "智能工单分派", "情绪识别与升级", "知识库自动更新"]),
    ("📝", "内容运营", ACCENT, "产出效率↑5倍",
     ["文章自动生成", "社交媒体发布", "数据报表编写", "多语言翻译"]),
    ("📈", "数据分析", PURPLE, "分析效率↑10倍",
     ["数据采集与清洗", "自动报表生成", "异常检测预警", "趋势预测分析"]),
    ("🔄", "流程自动化", YELLOW, "人工干预↓70%",
     ["审批流程处理", "合同审查比对", "财务对账结算", "系统间数据同步"]),
]

for i, (icon, title, color, metric, tasks) in enumerate(use_cases):
    col = i % 2
    row = i // 2
    left = Inches(0.8 + col * 6)
    top = Inches(1.2 + row * 3.1)

    card = add_rounded_rect(slide, left, top, Inches(5.6), Inches(2.8), BG_CARD)
    card.line.color.rgb = RGBColor(0x1A, 0x3A, 0x6A)
    card.line.width = Pt(1)

    # Header
    add_textbox(slide, left + Inches(0.3), top + Inches(0.2), Inches(3), Inches(0.5),
                f"{icon} {title}", font_size=26, color=WHITE, bold=True)

    # Metric badge
    badge = add_rounded_rect(slide, left + Inches(4.0), top + Inches(0.2), Inches(1.4), Inches(0.4),
                             RGBColor(0x0F, 0x28, 0x47))
    badge.line.color.rgb = color
    badge.line.width = Pt(1)
    add_textbox(slide, left + Inches(4.0), top + Inches(0.2), Inches(1.4), Inches(0.4),
                metric, font_size=16, color=color, bold=True, alignment=PP_ALIGN.CENTER)

    # Tasks
    for j, task in enumerate(tasks):
        add_circle(slide, left + Inches(0.3), top + Inches(0.85 + j * 0.45), Inches(0.18), color)
        add_textbox(slide, left + Inches(0.6), top + Inches(0.8 + j * 0.45), Inches(4.5), Inches(0.4),
                    task, font_size=18, color=MUTED)

# ════════════════════════════════════════════════════════════
# SLIDE 13: SECTION DIVIDER - Chapter 05
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, RGBColor(0x0F, 0x28, 0x47))

add_textbox(slide, Inches(2), Inches(2.2), Inches(9.333), Inches(0.8),
            "CHAPTER 05", font_size=24, color=ACCENT, bold=True,
            alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1.5), Inches(3.2), Inches(10.333), Inches(1.5),
            "成本与收益", font_size=56, color=WHITE, bold=True,
            alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(2), Inches(4.8), Inches(9.333), Inches(0.8),
            "ROI分析与投资回报", font_size=26, color=MUTED,
            alignment=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════
# SLIDE 14: ROI Analysis
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)

add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11.5), Inches(0.8),
            "05 投入产出分析", font_size=40, color=WHITE, bold=True)

# Cost section
add_textbox(slide, Inches(0.8), Inches(1.2), Inches(5.5), Inches(0.6),
            "投入成本", font_size=30, color=YELLOW, bold=True)

costs = [
    ("AI模型与平台", "¥20-50万/年"),
    ("开发集成费用", "¥10-30万/场景"),
    ("运维与培训", "¥5-10万/年"),
    ("合计（首年）", "¥50-100万"),
]
for i, (item, cost) in enumerate(costs):
    top = Inches(1.9 + i * 0.7)
    add_textbox(slide, Inches(0.8), top, Inches(5.5), Inches(0.5),
                item, font_size=22, color=MUTED)
    add_textbox(slide, Inches(4.5), top, Inches(2), Inches(0.5),
                cost, font_size=22, color=WHITE, bold=True)

# Return section
add_textbox(slide, Inches(7), Inches(1.2), Inches(5.5), Inches(0.6),
            "预期收益", font_size=30, color=GREEN, bold=True)

returns = [
    ("人力成本节省", "40-60%"),
    ("处理效率提升", "5-10倍"),
    ("错误率降低", "90%+"),
    ("投资回收期", "6-12个月"),
]
for i, (item, val) in enumerate(returns):
    top = Inches(1.9 + i * 0.7)
    add_textbox(slide, Inches(7), top, Inches(3.5), Inches(0.5),
                item, font_size=22, color=MUTED)
    add_textbox(slide, Inches(10.5), top, Inches(2), Inches(0.5),
                val, font_size=22, color=GREEN, bold=True)

# Conclusion box
conc = add_rounded_rect(slide, Inches(0.8), Inches(5.0), Inches(11.7), Inches(1.8),
                        RGBColor(0x0A, 0x2E, 0x1F))
conc.line.color.rgb = GREEN
conc.line.width = Pt(2)

add_multi_text(slide, Inches(1.2), Inches(5.2), Inches(10.9), Inches(1.4), [
    ("💡 结论", 28, GREEN, True, PP_ALIGN.CENTER),
    ("数字员工投入产出比显著，首年即可收回成本，长期ROI持续增长", 24, WHITE, False, PP_ALIGN.CENTER),
])

# ════════════════════════════════════════════════════════════
# SLIDE 15: SECTION DIVIDER - Chapter 06
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, RGBColor(0x0F, 0x28, 0x47))

add_textbox(slide, Inches(2), Inches(2.2), Inches(9.333), Inches(0.8),
            "CHAPTER 06", font_size=24, color=ACCENT, bold=True,
            alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1.5), Inches(3.2), Inches(10.333), Inches(1.5),
            "下一步计划", font_size=56, color=WHITE, bold=True,
            alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(2), Inches(4.8), Inches(9.333), Inches(0.8),
            "实施路线图与里程碑", font_size=26, color=MUTED,
            alignment=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════
# SLIDE 16: Next Steps
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)

add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11.5), Inches(0.8),
            "06 立即行动计划", font_size=40, color=WHITE, bold=True)

steps = [
    ("第1-2周", "需求调研与场景筛选", "访谈业务部门，识别高频重复性工作任务", "场景需求清单"),
    ("第3-4周", "技术方案设计与PoC验证", "完成架构设计，选取1个场景进行概念验证", "PoC验证报告"),
    ("第5-8周", "试点上线与迭代优化", "部署首个数字员工，收集反馈并持续优化", "试点运行数据"),
    ("第9-12周", "效果评估与规模推广", "总结试点经验，制定推广计划，扩大应用场景", "推广实施方案"),
]

for i, (week, action, detail, deliverable) in enumerate(steps):
    top = Inches(1.2 + i * 1.5)

    card = add_rounded_rect(slide, Inches(0.8), top, Inches(11.7), Inches(1.3), BG_CARD)
    card.line.color.rgb = RGBColor(0x1A, 0x3A, 0x6A)
    card.line.width = Pt(1)

    # Week badge
    wbadge = add_rounded_rect(slide, Inches(1.0), top + Inches(0.2), Inches(1.5), Inches(0.9),
                              RGBColor(0x0F, 0x28, 0x47))
    wbadge.line.color.rgb = ACCENT
    wbadge.line.width = Pt(2)
    add_textbox(slide, Inches(1.0), top + Inches(0.35), Inches(1.5), Inches(0.6),
                week, font_size=20, color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER)

    # Action & detail
    add_textbox(slide, Inches(2.8), top + Inches(0.15), Inches(6), Inches(0.5),
                action, font_size=24, color=WHITE, bold=True)
    add_textbox(slide, Inches(2.8), top + Inches(0.7), Inches(6), Inches(0.5),
                detail, font_size=18, color=MUTED)

    # Deliverable
    dbadge = add_rounded_rect(slide, Inches(10.0), top + Inches(0.3), Inches(2.2), Inches(0.7),
                              RGBColor(0x0A, 0x2E, 0x1F))
    dbadge.line.color.rgb = GREEN
    dbadge.line.width = Pt(1)
    add_textbox(slide, Inches(10.0), top + Inches(0.35), Inches(2.2), Inches(0.6),
                f"📋 {deliverable}", font_size=16, color=GREEN, bold=True,
                alignment=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════
# SLIDE 17: CLOSING
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)
add_shape_bg(slide, Inches(0), Inches(0), SLIDE_WIDTH, SLIDE_HEIGHT, BG_DARK)
add_circle(slide, Inches(9), Inches(-0.5), Inches(4), RGBColor(0x0F, 0x28, 0x47))
add_circle(slide, Inches(-1), Inches(5.5), Inches(3), RGBColor(0x0F, 0x28, 0x47))

add_textbox(slide, Inches(2), Inches(1.5), Inches(9.333), Inches(0.8),
            "THANK YOU", font_size=24, color=ACCENT, bold=True,
            alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1.5), Inches(2.5), Inches(10.333), Inches(1.5),
            "感谢聆听", font_size=52, color=ACCENT_LIGHT, bold=True,
            alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(2), Inches(4.2), Inches(9.333), Inches(1.0),
            "拥抱AI时代，让数字员工成为企业新生产力", font_size=26,
            color=MUTED, alignment=PP_ALIGN.CENTER)

add_shape_bg(slide, Inches(3), Inches(5.5), Inches(7.333), Inches(0.03), ACCENT)
add_textbox(slide, Inches(2), Inches(5.7), Inches(9.333), Inches(0.8),
            "敬请批评指正", font_size=22, color=MUTED, alignment=PP_ALIGN.CENTER)

# ─── Save ───
output_path = "/home/admin/.openclaw/workspace/digital-employee-report/数字员工落地汇报.pptx"
prs.save(output_path)
print(f"✅ PPT saved to: {output_path}")
