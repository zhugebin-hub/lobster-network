#!/usr/bin/env python3
"""
从图片创建PPT背景
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import shutil

# 源图片路径
source_image = "/home/admin/.openclaw/media/inbound/9ef2ef31-c17e-40fe-863a-b538b6b5592a.jpg"
# 输出PPT路径
output_pptx = "/home/admin/.openclaw/workspace/老师您费心了_PPT背景.pptx"

# 创建PPT
prs = Presentation()
prs.slide_width = Inches(13.333)  # 16:9 比例
prs.slide_height = Inches(7.5)

# 添加空白幻灯片
slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局

# 添加背景图片
left = top = Inches(0)
slide.shapes.add_picture(source_image, left, top, 
                         width=prs.slide_width, 
                         height=prs.slide_height)

# 添加文字内容（覆盖在背景上，增强可读性）
# 主标题
title_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(11.333), Inches(1.5))
title_frame = title_box.text_frame
title_para = title_frame.paragraphs[0]
title_para.text = "老师您费心了"
title_para.font.size = Pt(54)
title_para.font.bold = True
title_para.font.color.rgb = RGBColor(255, 255, 255)
title_para.alignment = PP_ALIGN.CENTER
title_frame.word_wrap = True

# 添加半透明背景框增强文字可读性
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR

# 五种心态内容
content_text = """
【释然型】
孩子已经尽力了，接受他的普通，平安健康就好

【心累型】
管不动了，随他去吧，我也需要放过自己

【保护型】
孩子不容易，不想再给他增加压力了

【防备型】
老师别总找我了，我们自己也头疼

【客套型】
嘴上说费心了，回家该说还说
"""

content_box = slide.shapes.add_textbox(Inches(1.5), Inches(2.8), Inches(10.333), Inches(4))
content_frame = content_box.text_frame
content_frame.word_wrap = True

# 设置内容文字
content_para = content_frame.paragraphs[0]
content_para.text = content_text.strip()
content_para.font.size = Pt(24)
content_para.font.color.rgb = RGBColor(255, 255, 255)
content_para.alignment = PP_ALIGN.LEFT

# 底部金句
footer_box = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(11.333), Inches(0.8))
footer_frame = footer_box.text_frame
footer_para = footer_frame.paragraphs[0]
footer_para.text = "今天我不为「说服」，只为交心"
footer_para.font.size = Pt(28)
footer_para.font.bold = True
footer_para.font.color.rgb = RGBColor(255, 215, 0)  # 金色
footer_para.alignment = PP_ALIGN.CENTER

# 保存PPT
prs.save(output_pptx)
print(f"PPT已生成：{output_pptx}")
