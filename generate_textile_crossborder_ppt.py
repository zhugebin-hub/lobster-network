#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
中职纺织品跨境电商教学 PPT 生成器
"""

from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# 创建演示文稿
prs = Presentation()
prs.slide_width = Cm(33.867)
prs.slide_height = Cm(19.05)

# 纺织品跨境电商主题颜色
COLORS = {
    'primary': (52, 73, 94),          # 深蓝灰
    'secondary': (155, 89, 182),      # 紫色
    'accent': (230, 126, 34),         # 橙色（强调）
    'success': (39, 174, 96),         # 绿色
    'light': (236, 240, 241),         # 浅灰
    'white': (255, 255, 255),         # 白色
    'dark': (20, 20, 30),             # 深色
    'gray': (127, 140, 141),          # 灰色
    'textile_blue': (41, 128, 185),   # 纺织蓝
}

def set_background(slide, color=COLORS['white']):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*color)

def add_header(slide, title, subtitle=""):
    # 顶部色带
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Cm(0), Cm(0), prs.slide_width, Cm(1.8)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*COLORS['primary'])
    shape.line.fill.background()
    
    # 装饰线（纺织纹理感）
    for i in range(5):
        line_x = Cm(22 + i * 2)
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, line_x, Cm(0.2), Cm(0.05), Cm(1.4)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*COLORS['accent'])
        shape.line.fill.background()
    
    # 标题
    textbox = slide.shapes.add_textbox(Cm(1.5), Cm(0.4), Cm(18), Cm(1))
    tf = textbox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = RGBColor(*COLORS['white'])
    p.font.name = 'Microsoft YaHei'
    
    if subtitle:
        textbox = slide.shapes.add_textbox(Cm(22), Cm(0.5), Cm(10), Cm(0.8))
        tf = textbox.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(*COLORS['white'])
        p.font.name = 'Microsoft YaHei'
        p.alignment = PP_ALIGN.RIGHT

def add_footer(slide, page_num=""):
    # 底线
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Cm(2), prs.slide_height - Cm(1), prs.slide_width - Cm(4), Cm(0.05)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*COLORS['textile_blue'])
    shape.line.fill.background()
    
    # 页码
    textbox = slide.shapes.add_textbox(prs.slide_width - Cm(3), prs.slide_height - Cm(1.3), Cm(2), Cm(0.5))
    tf = textbox.text_frame
    p = tf.paragraphs[0]
    p.text = page_num
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(*COLORS['gray'])
    p.font.name = 'Microsoft YaHei'
    p.alignment = PP_ALIGN.RIGHT

def add_title_slide(prs, title, subtitle, footer=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    
    # 背景装饰
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Cm(0), Cm(0), prs.slide_width, Cm(5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*COLORS['primary'])
    shape.line.fill.background()
    
    # 纺织纹理装饰线
    for i in range(15):
        y = Cm(1 + i * 0.25)
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Cm(0), y, prs.slide_width, Cm(0.03)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*COLORS['accent'])
        shape.fill.transparency = 0.5
        shape.line.fill.background()
    
    # 主标题
    textbox = slide.shapes.add_textbox(Cm(2), Cm(1.5), Cm(28), Cm(2.5))
    tf = textbox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = RGBColor(*COLORS['white'])
    p.font.name = 'Microsoft YaHei'
    
    # 副标题
    textbox = slide.shapes.add_textbox(Cm(2), Cm(4.5), Cm(20), Cm(1.5))
    tf = textbox.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(*COLORS['white'])
    p.font.name = 'Microsoft YaHei'
    
    # 底部信息
    if footer:
        textbox = slide.shapes.add_textbox(Cm(2), prs.slide_height - Cm(3), Cm(15), Cm(1))
        tf = textbox.text_frame
        p = tf.paragraphs[0]
        p.text = footer
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(*COLORS['gray'])
        p.font.name = 'Microsoft YaHei'
    
    return slide

def add_content_slide(prs, title, content_items, page_num=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_header(slide, title)
    
    y = Cm(3)
    for item in content_items:
        if isinstance(item, dict):
            if item.get('type') == 'title':
                textbox = slide.shapes.add_textbox(Cm(2), y, Cm(30), Cm(0.8))
                tf = textbox.text_frame
                p = tf.paragraphs[0]
                p.text = item['text']
                p.font.size = Pt(22)
                p.font.bold = True
                p.font.color.rgb = RGBColor(*COLORS['textile_blue'])
                p.font.name = 'Microsoft YaHei'
                y += Cm(1)
            elif item.get('type') == 'bullet':
                textbox = slide.shapes.add_textbox(Cm(2.5), y, Cm(29), Cm(0.7))
                tf = textbox.text_frame
                p = tf.paragraphs[0]
                p.text = "• " + item['text']
                p.font.size = Pt(18)
                p.font.color.rgb = RGBColor(*COLORS['dark'])
                p.font.name = 'Microsoft YaHei'
                y += Cm(0.9)
            elif item.get('type') == 'highlight':
                shape = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE, Cm(2), y, Cm(30), Cm(1.5)
                )
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor(*COLORS['light'])
                shape.line.color.rgb = RGBColor(*COLORS['textile_blue'])
                shape.line.width = Pt(2)
                
                textbox = slide.shapes.add_textbox(Cm(2.5), y + Cm(0.3), Cm(29), Cm(0.9))
                tf = textbox.text_frame
                p = tf.paragraphs[0]
                p.text = item['text']
                p.font.size = Pt(18)
                p.font.bold = True
                p.font.color.rgb = RGBColor(*COLORS['textile_blue'])
                p.font.name = 'Microsoft YaHei'
                y += Cm(2)
            elif item.get('type') == 'box':
                shape = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE, Cm(2), y, Cm(30), item.get('height', Cm(2))
                )
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor(*COLORS.get(item.get('color', 'light'), COLORS['light']))
                shape.line.fill.background()
                
                textbox = slide.shapes.add_textbox(Cm(2.5), y + Cm(0.3), Cm(29), Cm(1.5))
                tf = textbox.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = item['text']
                p.font.size = Pt(16)
                p.font.color.rgb = RGBColor(*COLORS['dark'])
                p.font.name = 'Microsoft YaHei'
                y += item.get('height', Cm(2)) + Cm(0.3)
        else:
            textbox = slide.shapes.add_textbox(Cm(2), y, Cm(30), Cm(0.7))
            tf = textbox.text_frame
            p = tf.paragraphs[0]
            p.text = item
            p.font.size = Pt(18)
            p.font.color.rgb = RGBColor(*COLORS['dark'])
            p.font.name = 'Microsoft YaHei'
            y += Cm(0.9)
    
    add_footer(slide, page_num)
    return slide

def add_platform_slide(prs, title, platforms, page_num=""):
    """平台介绍页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_header(slide, title)
    
    # 平台卡片
    card_width = Cm(14)
    card_height = Cm(5)
    start_x = Cm(2.5)
    start_y = Cm(3.5)
    
    for i, platform in enumerate(platforms):
        col = i % 2
        row = i // 2
        x = start_x + col * (card_width + Cm(1))
        y = start_y + row * (card_height + Cm(0.5))
        
        # 卡片背景
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, y, card_width, card_height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*COLORS['light'])
        shape.line.color.rgb = RGBColor(*COLORS['textile_blue'])
        shape.line.width = Pt(2)
        
        # 平台名称
        textbox = slide.shapes.add_textbox(x + Cm(1), y + Cm(0.5), card_width - Cm(2), Cm(0.8))
        tf = textbox.text_frame
        p = tf.paragraphs[0]
        p.text = platform['name']
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = RGBColor(*COLORS['textile_blue'])
        p.font.name = 'Microsoft YaHei'
        
        # 平台特点
        textbox = slide.shapes.add_textbox(x + Cm(1), y + Cm(1.5), card_width - Cm(2), Cm(3))
        tf = textbox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = platform['features']
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(*COLORS['dark'])
        p.font.name = 'Microsoft YaHei'
    
    add_footer(slide, page_num)
    return slide

def add_process_slide(prs, title, steps, page_num=""):
    """流程图页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_header(slide, title)
    
    # 流程步骤
    step_width = Cm(6)
    step_height = Cm(2.5)
    start_x = Cm(2)
    y = Cm(5)
    
    for i, step in enumerate(steps):
        x = start_x + i * (step_width + Cm(0.5))
        
        # 步骤框
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, y, step_width, step_height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*COLORS['textile_blue'])
        shape.line.fill.background()
        
        # 步骤文字
        textbox = slide.shapes.add_textbox(x + Cm(0.3), y + Cm(0.5), step_width - Cm(0.6), step_height - Cm(1))
        tf = textbox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = step
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = RGBColor(*COLORS['white'])
        p.font.name = 'Microsoft YaHei'
        p.alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        # 箭头
        if i < len(steps) - 1:
            arrow_x = x + step_width
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW, arrow_x, y + Cm(1), Cm(0.5), Cm(0.5)
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(*COLORS['accent'])
            shape.line.fill.background()
    
    add_footer(slide, page_num)
    return slide

def add_comparison_slide(prs, title, headers, rows, page_num=""):
    """对比表格页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_header(slide, title)
    
    col_width = Cm(6.5)
    row_height = Cm(1.5)
    start_x = Cm(2)
    start_y = Cm(3)
    
    # 表头
    for i, header in enumerate(headers):
        x = start_x + i * col_width
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, x, start_y, col_width, Cm(1.2)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*COLORS['textile_blue'])
        shape.line.fill.background()
        
        textbox = slide.shapes.add_textbox(x + Cm(0.3), start_y + Cm(0.2), col_width - Cm(0.6), Cm(0.8))
        tf = textbox.text_frame
        p = tf.paragraphs[0]
        p.text = header
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = RGBColor(*COLORS['white'])
        p.font.name = 'Microsoft YaHei'
        p.alignment = PP_ALIGN.CENTER
    
    # 内容行
    for row_idx, row_data in enumerate(rows):
        y = start_y + Cm(1.2) + row_idx * row_height
        for col_idx, cell in enumerate(row_data):
            x = start_x + col_idx * col_width
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, x, y, col_width - Cm(0.1), row_height - Cm(0.1)
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(*COLORS['light'])
            shape.line.color.rgb = RGBColor(*COLORS['gray'])
            shape.line.width = Pt(1)
            
            textbox = slide.shapes.add_textbox(x + Cm(0.3), y + Cm(0.3), col_width - Cm(0.9), row_height - Cm(0.6))
            tf = textbox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = cell
            p.font.size = Pt(12)
            p.font.color.rgb = RGBColor(*COLORS['dark'])
            p.font.name = 'Microsoft YaHei'
    
    add_footer(slide, page_num)
    return slide

# ============ 生成幻灯片 ============

# 1. 封面
add_title_slide(prs,
    title="纺织品跨境电商实务",
    subtitle="中职电子商务专业核心课程",
    footer="绍兴柯桥职校 • 2026 年 4 月"
)

# 2. 学习目标
add_content_slide(prs, "学习目标", [
    {"type": "highlight", "text": "学完本课程，你将能够："},
    "",
    {"type": "bullet", "text": "了解纺织品跨境电商的发展现状与趋势"},
    {"type": "bullet", "text": "熟悉主流跨境电商平台的特点与规则"},
    {"type": "bullet", "text": "掌握纺织品选品与Listing 优化技巧"},
    {"type": "bullet", "text": "学会跨境电商物流与支付操作"},
    {"type": "bullet", "text": "能够处理跨境售后与客户沟通"},
    {"type": "bullet", "text": "具备纺织品跨境电商创业基础能力"},
], "1/16")

# 3. 行业概述
add_content_slide(prs, "纺织品跨境电商行业概述", [
    {"type": "title", "text": "行业发展现状"},
    {"type": "bullet", "text": "中国纺织品出口额连续多年居全球第一"},
    {"type": "bullet", "text": "跨境电商成为纺织品出口新渠道"},
    {"type": "bullet", "text": "柯桥轻纺城年交易额超 3000 亿元"},
    {"type": "bullet", "text": "跨境电商 B2B、B2C 模式并行发展"},
    "",
    {"type": "title", "text": "主要优势"},
    {"type": "bullet", "text": "产业链完整，货源充足"},
    {"type": "bullet", "text": "性价比高，国际竞争力强"},
    {"type": "bullet", "text": "柔性供应链，小单快反能力"},
    {"type": "bullet", "text": "电商基础设施完善"},
], "2/16")

# 4. 纺织品分类
add_content_slide(prs, "纺织品跨境电商产品分类", [
    {"type": "title", "text": "按用途分类"},
    {"type": "bullet", "text": "服装类：T 恤、衬衫、连衣裙、外套等"},
    {"type": "bullet", "text": "家纺类：床品四件套、窗帘、毛巾等"},
    {"type": "bullet", "text": "面料类：棉布、丝绸、化纤面料等"},
    {"type": "bullet", "text": "辅料类：拉链、纽扣、蕾丝花边等"},
    "",
    {"type": "title", "text": "按市场分类"},
    {"type": "bullet", "text": "欧美市场：注重品质、环保认证"},
    {"type": "bullet", "text": "东南亚市场：价格敏感、中低端为主"},
    {"type": "bullet", "text": "中东市场：偏好鲜艳色彩、传统图案"},
    {"type": "bullet", "text": "日韩市场：注重设计、细节精致"},
], "3/16")

# 5. 主流平台
add_platform_slide(prs, "主流跨境电商平台", [
    {"name": "阿里巴巴国际站", "features": "• B2B 平台，适合批量订单\n• 柯桥商家入驻多\n• 支持 RFQ 报价\n• 信保订单保障"},
    {"name": "Amazon", "features": "• 全球最大 B2C 平台\n• 流量大、客单价高\n• FBA 物流体系完善\n• 规则严格，需合规运营"},
    {"name": "SHEIN", "features": "• 快时尚跨境电商\n• 小单快反模式\n• 适合服装品类\n• 供应商合作模式"},
    {"name": "速卖通", "features": "• 阿里旗下 B2C 平台\n• 覆盖 200+ 国家\n• 适合中小卖家\n• 运营门槛相对较低"},
], "4/16")

# 6. 平台对比
add_comparison_slide(prs, "平台特点对比",
    ["平台", "模式", "主要市场", "入驻费用", "适合卖家"],
    [
        ["阿里巴巴国际站", "B2B", "全球", "2.98 万/年起", "工厂/贸易商"],
        ["Amazon", "B2C", "欧美为主", "约 3 万/年", "品牌卖家"],
        ["SHEIN", "B2C", "全球", "供应商合作", "服装工厂"],
        ["速卖通", "B2C", "新兴市场", "1 万/年", "中小卖家"],
        ["TikTok Shop", "社交电商", "年轻用户", "较低", "内容创作者"],
    ],
    "5/16"
)

# 7. 选品策略
add_content_slide(prs, "纺织品选品策略", [
    {"type": "highlight", "text": "选品是跨境电商成功的关键！"},
    "",
    {"type": "title", "text": "选品原则"},
    {"type": "bullet", "text": "市场需求大：搜索量高、竞争适中"},
    {"type": "bullet", "text": "利润空间足：毛利率 30% 以上"},
    {"type": "bullet", "text": "物流友好：重量轻、不易损坏"},
    {"type": "bullet", "text": "合规风险低：符合目标国标准"},
    "",
    {"type": "title", "text": "选品工具"},
    {"type": "bullet", "text": "Google Trends 趋势分析"},
    {"type": "bullet", "text": "平台热销榜、关键词工具"},
    {"type": "bullet", "text": "第三方数据工具（卖家精灵、Jungle Scout）"},
], "6/16")

# 8. Listing 优化
add_content_slide(prs, "商品 Listing 优化", [
    {"type": "title", "text": "标题优化"},
    {"type": "bullet", "text": "核心关键词前置"},
    {"type": "bullet", "text": "包含材质、款式、颜色、尺码"},
    {"type": "bullet", "text": "例：100% Cotton T-Shirt Men Casual Short Sleeve Summer Breathable"},
    "",
    {"type": "title", "text": "图片优化"},
    {"type": "bullet", "text": "主图白底清晰，展示产品全貌"},
    {"type": "bullet", "text": "细节图展示面料纹理、做工"},
    {"type": "bullet", "text": "场景图展示穿着/使用效果"},
    {"type": "bullet", "text": "尺寸图标注详细规格"},
    "",
    {"type": "title", "text": "描述优化"},
    {"type": "bullet", "text": "突出卖点和差异化"},
    {"type": "bullet", "text": "面料成分、洗涤说明"},
    {"type": "bullet", "text": "尺码表（欧美/亚洲码对照）"},
], "7/16")

# 9. 跨境物流
add_process_slide(prs, "跨境电商物流流程", [
    "备货入仓",
    "打包贴单",
    "报关出口",
    "国际运输",
    "清关入境",
    "末端配送",
    "客户签收"
], "8/16")

# 10. 物流方式对比
add_comparison_slide(prs, "物流方式对比",
    ["物流方式", "时效", "价格", "适合产品", "追踪"],
    [
        ["国际快递", "3-7 天", "高", "样品/急件", "全程追踪"],
        ["空运专线", "7-15 天", "中高", "高价值产品", "全程追踪"],
        ["海运专线", "25-45 天", "低", "大宗货物", "部分追踪"],
        ["中欧班列", "15-25 天", "中", "欧洲线路", "全程追踪"],
        ["海外仓", "2-5 天", "中", "热销爆款", "本地追踪"],
    ],
    "9/16"
)

# 11. 跨境支付
add_content_slide(prs, "跨境支付与结算", [
    {"type": "title", "text": "主流收款方式"},
    {"type": "bullet", "text": "平台收款：Amazon Pay、Alipay、PayPal"},
    {"type": "bullet", "text": "第三方收款：PingPong、连连支付、空中云汇"},
    {"type": "bullet", "text": "银行电汇：T/T、L/C（B2B 常用）"},
    "",
    {"type": "title", "text": "结汇流程"},
    {"type": "bullet", "text": "外币收款 → 第三方支付账户 → 人民币结汇 → 提现到银行卡"},
    "",
    {"type": "title", "text": "注意事项"},
    {"type": "bullet", "text": "关注汇率波动，适时结汇"},
    {"type": "bullet", "text": "了解各平台手续费"},
    {"type": "bullet", "text": "合规申报，依法纳税"},
], "10/16")

# 12. 客户服务
add_content_slide(prs, "跨境客户服务与沟通", [
    {"type": "title", "text": "售前咨询"},
    {"type": "bullet", "text": "及时回复询盘（24 小时内）"},
    {"type": "bullet", "text": "专业解答产品问题（面料、尺码、颜色）"},
    {"type": "bullet", "text": "提供搭配建议、定制方案"},
    "",
    {"type": "title", "text": "售后处理"},
    {"type": "bullet", "text": "质量问题：退换货/部分退款"},
    {"type": "bullet", "text": "物流问题：协助查询、补发"},
    {"type": "bullet", "text": "尺寸问题：提供换货或优惠券"},
    "",
    {"type": "title", "text": "沟通技巧"},
    {"type": "bullet", "text": "使用简洁清晰的英语"},
    {"type": "bullet", "text": "保持礼貌专业"},
    {"type": "bullet", "text": "善用模板提高效率"},
], "11/16")

# 13. 常见纠纷
add_content_slide(prs, "常见纠纷与处理", [
    {"type": "title", "text": "质量问题纠纷"},
    {"type": "bullet", "text": "原因：色差、面料不符、做工瑕疵"},
    {"type": "bullet", "text": "处理：道歉 + 退换货/部分退款 + 改进"},
    "",
    {"type": "title", "text": "物流纠纷"},
    {"type": "bullet", "text": "原因：延迟、丢包、破损"},
    {"type": "bullet", "text": "处理：协助查询 + 补发/退款 + 物流索赔"},
    "",
    {"type": "title", "text": "预防策略"},
    {"type": "bullet", "text": "如实描述，避免过度承诺"},
    {"type": "bullet", "text": "严格质检，把控出货质量"},
    {"type": "bullet", "text": "购买物流保险，降低风险"},
    {"type": "bullet", "text": "保留证据，便于纠纷处理"},
], "12/16")

# 14. 合规要求
add_content_slide(prs, "纺织品出口合规要求", [
    {"type": "highlight", "text": "合规是跨境电商的生命线！"},
    "",
    {"type": "title", "text": "产品认证"},
    {"type": "bullet", "text": "欧盟：CE 认证、OEKO-TEX 标准 100"},
    {"type": "bullet", "text": "美国：CPC 认证（儿童产品）、CA 65"},
    {"type": "bullet", "text": "日本：JIS 标准、PSE 认证"},
    "",
    {"type": "title", "text": "标签要求"},
    {"type": "bullet", "text": "成分标签（纤维含量）"},
    {"type": "bullet", "text": "洗涤标识"},
    {"type": "bullet", "text": "原产地标签（Made in China）"},
    "",
    {"type": "title", "text": "知识产权"},
    {"type": "bullet", "text": "不侵犯商标、专利、版权"},
    {"type": "bullet", "text": "自主设计，避免仿款"},
], "13/16")

# 15. 实战案例
add_content_slide(prs, "柯桥纺织品跨境电商案例", [
    {"type": "title", "text": "案例一：某家纺企业 Amazon 运营"},
    {"type": "bullet", "text": "产品：床品四件套"},
    {"type": "bullet", "text": "策略：FBA 发货 + 精品运营"},
    {"type": "bullet", "text": "成果：月销 5000 单，毛利率 35%"},
    "",
    {"type": "title", "text": "案例二：某面料厂阿里巴巴国际站"},
    {"type": "bullet", "text": "产品：功能性面料"},
    {"type": "bullet", "text": "策略：RFQ 报价 + 直播看厂"},
    {"type": "bullet", "text": "成果：年出口额 500 万美元"},
    "",
    {"type": "title", "text": "案例三：SHEIN 供应商合作"},
    {"type": "bullet", "text": "产品：快时尚女装"},
    {"type": "bullet", "text": "策略：小单快反 + 柔性供应"},
    {"type": "bullet", "text": "成果：日均出货 10000 件"},
], "14/16")

# 16. 创业指导
add_content_slide(prs, "纺织品跨境电商创业指导", [
    {"type": "title", "text": "启动准备"},
    {"type": "bullet", "text": "资金：5-10 万起步（平台 + 备货 + 物流）"},
    {"type": "bullet", "text": "货源：柯桥轻纺城市场/1688"},
    {"type": "bullet", "text": "团队：运营 + 客服 + 美工（可兼职）"},
    "",
    {"type": "title", "text": "风险提示"},
    {"type": "bullet", "text": "库存风险：避免大量备货"},
    {"type": "bullet", "text": "汇率风险：关注汇率波动"},
    {"type": "bullet", "text": "平台风险：规则变化、封号风险"},
    "",
    {"type": "title", "text": "学习建议"},
    {"type": "bullet", "text": "参加平台官方培训"},
    {"type": "bullet", "text": "加入行业社群交流"},
    {"type": "bullet", "text": "持续学习，关注行业动态"},
], "15/16")

# 17. 课后练习
add_content_slide(prs, "课后练习与实训", [
    {"type": "highlight", "text": "实训任务："},
    "",
    {"type": "title", "text": "基础任务"},
    {"type": "bullet", "text": "调研 3 个跨境电商平台，完成对比表格"},
    {"type": "bullet", "text": "选择一款纺织品，撰写 Listing 文案"},
    {"type": "bullet", "text": "计算产品定价（含成本、运费、利润）"},
    "",
    {"type": "title", "text": "提高任务"},
    {"type": "bullet", "text": "模拟开通店铺，完成产品上架"},
    {"type": "bullet", "text": "设计营销方案（优惠券/促销活动）"},
    {"type": "bullet", "text": "编写客服话术模板（英文）"},
    "",
    {"type": "title", "text": "拓展资源"},
    {"type": "bullet", "text": "阿里巴巴国际站大学、Amazon 卖家大学"},
    {"type": "bullet", "text": "雨果网、跨境眼等行业媒体"},
], "16/16")

# 18. 结束页
add_title_slide(prs,
    title="感谢观看",
    subtitle="欢迎提问交流",
    footer="纺织品跨境电商实务 • 2026 年 4 月"
)

# 保存文件
output_path = '/home/admin/.openclaw/workspace/纺织品跨境电商教学 PPT.pptx'
prs.save(output_path)
print(f"PPT 已生成：{output_path}")
print(f"共 {len(prs.slides)} 页")
