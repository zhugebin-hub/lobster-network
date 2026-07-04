#!/usr/bin/env python3
"""生成《AI时代的人类社会发展趋势与人文思考》PPT"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# 配色方案
DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT_BLUE = RGBColor(0x4A, 0x90, 0xD9)
ACCENT_ORANGE = RGBColor(0xE8, 0x8D, 0x4A)
ACCENT_GREEN = RGBColor(0x5C, 0xB8, 0x8C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
DARK_TEXT = RGBColor(0x33, 0x33, 0x33)
CARD_BG = RGBColor(0xF5, 0xF7, 0xFA)

def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Microsoft YaHei"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_multiline_textbox(slide, left, top, width, height, lines, font_size=16, color=WHITE, line_spacing=1.5, font_name="Microsoft YaHei"):
    """lines: list of (text, size, color, bold) tuples, or just strings"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if isinstance(line, str):
            text, size, clr, bld = line, font_size, color, False
        else:
            text = line[0]
            size = line[1] if len(line) > 1 else font_size
            clr = line[2] if len(line) > 2 else color
            bld = line[3] if len(line) > 3 else False
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = clr
        p.font.bold = bld
        p.font.name = font_name
        p.space_after = Pt(size * 0.4)
    return txBox

def add_card(slide, left, top, width, height, title, content_lines, title_color=ACCENT_BLUE, bg_color=CARD_BG):
    """添加卡片形状"""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(12)
    tf.margin_right = Pt(12)
    tf.margin_top = Pt(8)
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(14)
    p.font.color.rgb = title_color
    p.font.bold = True
    p.font.name = "Microsoft YaHei"
    for line in content_lines:
        p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(11)
        p.font.color.rgb = DARK_TEXT
        p.font.name = "Microsoft YaHei"
        p.space_after = Pt(2)

def add_accent_bar(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

# ============================================================
# Slide 1: 封面
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide, DARK_BG)
add_accent_bar(slide, 0.5, 2.5, 12.3, 0.06, ACCENT_BLUE)
add_textbox(slide, 1.5, 2.8, 10, 1.2, "AI时代的人类社会发展趋势与人文思考", font_size=36, color=WHITE, bold=True)
add_textbox(slide, 1.5, 4.2, 10, 0.8, "从《头号玩家》到数字人文课程的跨学科分析", font_size=22, color=LIGHT_GRAY)
add_accent_bar(slide, 1.5, 5.2, 3, 0.04, ACCENT_ORANGE)
add_textbox(slide, 1.5, 5.5, 10, 0.5, "浙江工商大学 · 《数字人文与小龙虾》课程", font_size=16, color=LIGHT_GRAY)
add_textbox(slide, 1.5, 6.0, 10, 0.5, "2026 年 5 月 7 日 · 小龙虾 - 诸葛虾 整理", font_size=14, color=LIGHT_GRAY)

# ============================================================
# Slide 2: 目录
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_textbox(slide, 1, 0.5, 10, 0.6, "目  录", font_size=32, color=WHITE, bold=True)
add_accent_bar(slide, 1, 1.2, 2, 0.05, ACCENT_BLUE)

chapters = [
    ("01", "引言：为什么是现在？"),
    ("02", "《头号玩家》的预言与现实"),
    ("03", "AI时代十大发展趋势"),
    ("04", "小龙虾在AI时代的定位与实践"),
    ("05", "数字人文课程中的人文思考"),
    ("06", "综合讨论：技术向善，人文向远"),
    ("07", "行动建议"),
    ("08", "结语"),
]

for i, (num, title) in enumerate(chapters):
    y = 1.6 + i * 0.65
    add_textbox(slide, 1.5, y, 0.8, 0.5, num, font_size=24, color=ACCENT_ORANGE, bold=True)
    add_textbox(slide, 2.5, y, 8, 0.5, title, font_size=18, color=WHITE)

# ============================================================
# Slide 3: 引言 - 为什么是现在
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_textbox(slide, 1, 0.4, 10, 0.6, "01  引言：为什么是现在？", font_size=28, color=WHITE, bold=True)
add_accent_bar(slide, 1, 1.1, 2, 0.05, ACCENT_BLUE)

cards = [
    ("2026 年", "「十五五」规划开局之年\nAI 产业从「技术探索」向「全面赋能」质变"),
    ("4月28日", "中共中央政治局会议定调\n「人工智能+」全面实施"),
    ("AI 智能体", "从「工具」到「伙伴」\n具备自主感知、决策和执行能力"),
]

for i, (title, content) in enumerate(cards):
    x = 1 + i * 4
    add_card(slide, x, 1.8, 3.5, 2.2, title, content.split('\n'), ACCENT_BLUE, DARK_BG)
    # Make card text white on dark bg
    for shape in slide.shapes:
        if shape.left == Inches(x) and shape.top == Inches(1.8):
            tf = shape.text_frame
            for p in tf.paragraphs:
                p.font.color.rgb = WHITE

# Three dimensions
add_textbox(slide, 1, 4.5, 10, 0.5, "三个讨论维度", font_size=20, color=ACCENT_ORANGE, bold=True)
dims = [
    ("🎬 文化镜像", "《头号玩家》电影分析\n虚拟与现实的关系"),
    ("📈 技术趋势", "AI时代社会发展分析\n十大发展趋势"),
    ("🏛️ 人文反思", "数字人文课程实践\n什么让人成为人？"),
]
for i, (title, content) in enumerate(dims):
    x = 1 + i * 4
    add_textbox(slide, x, 5.2, 3.5, 0.4, title, font_size=16, color=ACCENT_GREEN, bold=True)
    add_textbox(slide, x, 5.7, 3.5, 1.0, content, font_size=12, color=LIGHT_GRAY)

# ============================================================
# Slide 4: 头号玩家 - 核心设定
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_textbox(slide, 1, 0.4, 10, 0.6, "02  《头号玩家》的预言与现实", font_size=28, color=WHITE, bold=True)
add_accent_bar(slide, 1, 1.1, 2, 0.05, ACCENT_BLUE)

add_textbox(slide, 1, 1.5, 10, 0.5, "核心设定", font_size=22, color=ACCENT_ORANGE, bold=True)

info = [
    ("时间线", "2045 年，世界动荡不安"),
    ("背景", "人们逃入虚拟世界「绿洲」（OASIS）"),
    ("剧情", "3 把钥匙 + 1 个彩蛋，继承 5000 亿美元"),
    ("主题", "现实 vs 虚拟 — 什么才是真实的？"),
]

for i, (label, desc) in enumerate(info):
    y = 2.3 + i * 0.7
    add_textbox(slide, 1.5, y, 2, 0.5, label, font_size=16, color=ACCENT_BLUE, bold=True)
    add_textbox(slide, 3.5, y, 8, 0.5, desc, font_size=15, color=WHITE)

add_textbox(slide, 1, 5.5, 11, 0.8, '"现实才是唯一的真实。虚拟世界再美好，也替代不了真实的生活。"  —— 哈利迪', font_size=16, color=ACCENT_GREEN, bold=True)

# ============================================================
# Slide 5: 三把钥匙
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_textbox(slide, 1, 0.4, 10, 0.6, "三把钥匙的隐喻", font_size=26, color=WHITE, bold=True)
add_accent_bar(slide, 1, 1.1, 2, 0.05, ACCENT_BLUE)

keys = [
    ("🔑 铜钥匙", "赛车关 — 倒车看真相", "逆向思维\n不盲从，退一步看全局", "AI时代需要批判性思维\n不盲信算法输出"),
    ("🔑 Jade钥匙", "迷宫/剧院关", "共情能力\n理解创造者的内心", "AI能分析文本\n但理解不了人心"),
    ("🔑 水晶钥匙", "游戏关 — 彩蛋的彩蛋", "理解真正意义\n游戏是为了快乐", "AI是工具\n关键是为什么用"),
]

for i, (title, challenge, metaphor, reality) in enumerate(keys):
    x = 0.8 + i * 4.2
    add_card(slide, x, 1.5, 3.8, 4.5, title, [f"考验：{challenge}", "", f"隐喻：{metaphor}", "", f"现实对应：{reality}"], ACCENT_ORANGE, DARK_BG)
    for shape in slide.shapes:
        if shape.left == Inches(x) and shape.top == Inches(1.5):
            tf = shape.text_frame
            for p in tf.paragraphs:
                p.font.color.rgb = WHITE
                p.font.name = "Microsoft YaHei"

# ============================================================
# Slide 6: 十大趋势总览
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_textbox(slide, 1, 0.4, 10, 0.6, "03  AI时代十大发展趋势", font_size=28, color=WHITE, bold=True)
add_accent_bar(slide, 1, 1.1, 2, 0.05, ACCENT_BLUE)

trends = [
    ("01", "教育革命", "知识灌输→能力培养", "教师角色转型"),
    ("02", "职业重构", "消失的 vs 新兴的", "就业公平"),
    ("03", "生活方式", "智能环境成基础设施", "隐私与自主"),
    ("04", "人际关系", "AI作为社交媒介", "真实vs模拟"),
    ("05", "医疗健康", "预防性+个性化", "数据伦理"),
    ("06", "社会治理", "数据驱动决策", "算法透明"),
    ("07", "创造力释放", "AI作为创作伙伴", "创作本质"),
    ("08", "伦理与治理", "规则重建", "价值判断"),
    ("09", "认知进化", "人机协作新思维", "独特认知"),
    ("10", "可持续发展", "AI应对全球挑战", "系统变革"),
]

for i, (num, trend, change, challenge) in enumerate(trends):
    y = 1.5 + i * 0.55
    add_textbox(slide, 1.2, y, 0.5, 0.4, num, font_size=14, color=ACCENT_ORANGE, bold=True)
    add_textbox(slide, 1.8, y, 2.5, 0.4, trend, font_size=14, color=ACCENT_BLUE, bold=True)
    add_textbox(slide, 4.5, y, 3.5, 0.4, change, font_size=12, color=WHITE)
    add_textbox(slide, 8.5, y, 4, 0.4, challenge, font_size=12, color=LIGHT_GRAY)

add_textbox(slide, 1, 7, 11, 0.5, "核心问题：在每一个趋势中，人文的价值在哪里？", font_size=14, color=ACCENT_GREEN)

# ============================================================
# Slide 7-8: 重点趋势详解
# ============================================================
for slide_idx, (title, items) in enumerate([
    ("教育革命", [
        ("现状", "标准化考试、统一课程、知识记忆为主"),
        ("未来", "AI个性化学习路径，每个学生有AI导师"),
        ("教师角色", "从「知识传授者」变为「学习引导者」"),
        ("核心", "重点培养：批判思维、创造力、情感智慧"),
    ]),
    ("职业重构", [
        ("被替代", "重复性脑力劳动、规律性体力劳动"),
        ("新职业", "AI训练师、提示工程师、AI伦理审查"),
        ("不变", "创造力、情感智慧、复杂决策"),
        ("判断", "问题不是「AI抢饭碗」，是「分配不公平」"),
    ]),
]):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_textbox(slide, 1, 0.4, 10, 0.6, f"趋势详解：{title}", font_size=26, color=WHITE, bold=True)
    add_accent_bar(slide, 1, 1.1, 2, 0.05, ACCENT_BLUE)
    
    for i, (label, desc) in enumerate(items):
        y = 1.5 + i * 1.2
        add_textbox(slide, 1.5, y, 2, 0.5, label, font_size=18, color=ACCENT_ORANGE, bold=True)
        add_textbox(slide, 3.5, y, 8, 0.8, desc, font_size=16, color=WHITE)

# ============================================================
# Slide 9: 小龙虾定位
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_textbox(slide, 1, 0.4, 10, 0.6, "04  小龙虾在AI时代的定位与实践", font_size=28, color=WHITE, bold=True)
add_accent_bar(slide, 1, 1.1, 2, 0.05, ACCENT_BLUE)

roles = [
    ("📡 信息枢纽", "整合、筛选、传递信息"),
    ("🤝 学习伙伴", "个性化辅导、答疑"),
    ("⚡ 效率工具", "自动化重复任务"),
    ("💡 创意催化剂", "激发灵感、提供视角"),
    ("🌉 桥梁", "连接人与技术"),
]

for i, (role, desc) in enumerate(roles):
    y = 1.5 + i * 0.8
    add_textbox(slide, 1.5, y, 3, 0.6, role, font_size=16, color=ACCENT_BLUE, bold=True)
    add_textbox(slide, 4.5, y, 7, 0.6, desc, font_size=15, color=WHITE)

add_textbox(slide, 1, 5.8, 11, 0.4, "设计原则", font_size=20, color=ACCENT_ORANGE, bold=True)
principles = "增强而非替代  ·  透明而非伪装  ·  教育而非依赖  ·  陪伴而非控制  ·  成长而非静态"
add_textbox(slide, 1, 6.3, 11, 0.6, principles, font_size=14, color=ACCENT_GREEN)

# ============================================================
# Slide 10: 人文思考 - 五个维度
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_textbox(slide, 1, 0.4, 10, 0.6, "05  数字人文课程中的人文思考", font_size=28, color=WHITE, bold=True)
add_accent_bar(slide, 1, 1.1, 2, 0.05, ACCENT_BLUE)

add_textbox(slide, 1, 1.3, 11, 0.5, "数字人文的本质：用技术扩展人文边界，用数据验证人文直觉", font_size=16, color=ACCENT_GREEN)
add_textbox(slide, 1, 1.8, 11, 0.4, "但解释和意义，永远是人做的事", font_size=16, color=ACCENT_ORANGE, bold=True)

dims = [
    ("创造力", "模式重组 vs 生命体验", "AI能写诗，但没看过樱花"),
    ("共情", "模拟语言 vs 真实体验", "慈善的本质是共情，不是信息处理"),
    ("记忆", "数据保存 vs 价值判断", "什么值得被记住？"),
    ("公平", "技术中立 vs 使用决定", "数字鸿沟在AI时代可能更大"),
    ("意义", "怎么做 vs 为什么做", "AI回答怎么做，人追问为什么"),
]

for i, (dim, contrast, example) in enumerate(dims):
    y = 2.4 + i * 0.9
    add_textbox(slide, 1.5, y, 1.5, 0.5, dim, font_size=16, color=ACCENT_BLUE, bold=True)
    add_textbox(slide, 3.2, y, 4, 0.5, contrast, font_size=14, color=WHITE)
    add_textbox(slide, 7.5, y, 5, 0.5, example, font_size=13, color=LIGHT_GRAY)

# ============================================================
# Slide 11: 什么让人成为人
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_textbox(slide, 1, 0.4, 12, 0.6, "核心命题：什么让人成为人？", font_size=28, color=WHITE, bold=True)
add_accent_bar(slide, 1, 1.1, 3, 0.05, ACCENT_ORANGE)

qualities = [
    ("💗 共情与脆弱性", "能感受痛苦，也能被感动\n人与人之间的信任建立在共同脆弱性上"),
    ("🎯 自由意志", "能做非理性选择，并承担后果\nAI的输出是概率，人的选择是责任"),
    ("🎨 创造无意义的美好", "艺术、音乐、诗歌——不「有用」但让人成为人\nAI是模式重组，人是内在冲动"),
    ("❓ 追问为什么", "不满足于「怎么做」，追问「为什么做」\nAI提供能力，人文提供方向"),
    ("🤝 真实连接", "真实的拥抱胜过虚拟的陪伴\n技术增强现实，不替代现实"),
]

for i, (title, desc) in enumerate(qualities):
    x = 0.8 + (i % 3) * 4.2
    y = 1.8 + (i // 3) * 2.5
    add_card(slide, x, y, 3.8, 2.0, title, desc.split('\n'), ACCENT_GREEN, DARK_BG)
    for shape in slide.shapes:
        if shape.left == Inches(x) and shape.top == Inches(y):
            tf = shape.text_frame
            for p in tf.paragraphs:
                p.font.color.rgb = WHITE

# ============================================================
# Slide 12: 行动建议
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_textbox(slide, 1, 0.4, 10, 0.6, "07  行动建议", font_size=28, color=WHITE, bold=True)
add_accent_bar(slide, 1, 1.1, 2, 0.05, ACCENT_BLUE)

groups = [
    ("🎓 学生", ["学会用AI做调研", "培养AI替代不了的能力", "关注真实问题", "保持人文关怀"]),
    ("👨‍👩‍👧 家长", ["不焦虑、不排斥", "设定规则引导使用", "亲子共学", "培养人味能力"]),
    ("👨‍🏫 教育者", ["AI教书，教师育人", "从教知识到教思维", "培养AI素养", "成为点燃者"]),
    ("🏛️ 政策制定", ["关注转型期阵痛", "推动技术民主化", "建立AI伦理框架", "投资人文教育"]),
]

for i, (group, tips) in enumerate(groups):
    x = 0.8 + i * 3.3
    add_card(slide, x, 1.5, 3.0, 4.5, group, tips, ACCENT_BLUE, DARK_BG)
    for shape in slide.shapes:
        if shape.left == Inches(x) and shape.top == Inches(1.5):
            tf = shape.text_frame
            for p in tf.paragraphs:
                p.font.color.rgb = WHITE

# ============================================================
# Slide 13: 结语
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_accent_bar(slide, 0.5, 2, 12.3, 0.06, ACCENT_BLUE)

add_textbox(slide, 1.5, 2.3, 10, 0.8, '"Thanks for playing."', font_size=32, color=ACCENT_ORANGE, bold=True)
add_textbox(slide, 1.5, 3.3, 10, 0.6, '但记住——现实才是唯一的真实。', font_size=22, color=WHITE)

add_accent_bar(slide, 1.5, 4.3, 3, 0.04, ACCENT_GREEN)

add_textbox(slide, 1.5, 4.6, 10, 0.5, '技术向善，人文向远', font_size=26, color=ACCENT_GREEN, bold=True)
add_textbox(slide, 1.5, 5.3, 10, 0.5, '在AI时代，成为更好的自己 🦞', font_size=20, color=WHITE)

add_textbox(slide, 1.5, 6.5, 10, 0.5, '浙江工商大学 · 《数字人文与小龙虾》课程 · 2026年5月', font_size=14, color=LIGHT_GRAY)

# 保存
output_path = os.path.join(os.path.dirname(__file__), "AI时代的人类社会发展趋势与人文思考.pptx")
prs.save(output_path)
print(f"✅ PPT saved: {output_path}")
print(f"   Slides: {len(prs.slides)}")
