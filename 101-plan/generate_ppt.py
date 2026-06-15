#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
计算机网络"101 计划"课程建设第一次研讨会汇报 PPT
生成时间：2026-06-01
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

def set_slide_bg(slide, r, g, b):
    """设置幻灯片背景色"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(r, g, b)

def add_title_slide(prs, title, subtitle):
    """添加标题幻灯片"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
    set_slide_bg(slide, 0x1B, 0x3A, 0x5C)  # 深蓝色背景
    
    # 主标题
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.alignment = PP_ALIGN.CENTER
    
    # 副标题
    txBox2 = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(8), Inches(1))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = subtitle
    p2.font.size = Pt(24)
    p2.font.color.rgb = RGBColor(0xA0, 0xC4, 0xE8)
    p2.alignment = PP_ALIGN.CENTER
    
    # 底部信息
    txBox3 = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(8), Inches(0.8))
    tf3 = txBox3.text_frame
    tf3.word_wrap = True
    p3 = tf3.paragraphs[0]
    p3.text = "2026 年 6 月 1 日"
    p3.font.size = Pt(18)
    p3.font.color.rgb = RGBColor(0x80, 0x80, 0x80)
    p3.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, content_list, is_table=False):
    """添加内容幻灯片"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, 0xFF, 0xFF, 0xFF)
    
    # 标题
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x1B, 0x3A, 0x5C)
    
    # 下划线
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 
        Inches(0.5), Inches(1.1), Inches(3), Inches(0.05)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x2E, 0x86, 0xC1)
    shape.line.fill.background()
    
    if is_table:
        # 添加表格
        rows = len(content_list)
        cols = len(content_list[0]) if content_list else 0
        table_shape = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.3), Inches(9), Inches(4.5))
        table = table_shape.table
        
        for i, row_data in enumerate(content_list):
            for j, cell_text in enumerate(row_data):
                cell = table.cell(i, j)
                cell.text = str(cell_text)
                
                # 设置单元格格式
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.size = Pt(11)
                    if i == 0:
                        paragraph.font.bold = True
                        paragraph.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                    else:
                        paragraph.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
                
                # 设置单元格背景
                if i == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(0x2E, 0x86, 0xC1)
                elif i % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(0xF0, 0xF8, 0xFF)
        
        return slide
    
    # 添加文本内容
    txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9), Inches(4.5))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    
    for i, item in enumerate(content_list):
        if i == 0:
            p = tf2.paragraphs[0]
        else:
            p = tf2.add_paragraph()
        
        p.text = item
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
        p.space_after = Pt(8)
        
        # 项目符号
        if item.startswith("•"):
            p.level = 1
            p.font.size = Pt(14)
    
    return slide

def add_milestone_slide(prs):
    """添加里程碑幻灯片"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, 0xFF, 0xFF, 0xFF)
    
    # 标题
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "关键里程碑"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x1B, 0x3A, 0x5C)
    
    # 下划线
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 
        Inches(0.5), Inches(1.1), Inches(3), Inches(0.05)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x2E, 0x86, 0xC1)
    shape.line.fill.background()
    
    # 里程碑数据
    milestones = [
        ("6 月", "知识体系初稿 + 教学大纲终稿", "诸葛斌"),
        ("7 月", "教材样章 2 章 + 实验案例 5 个", "高明/蒋献"),
        ("8 月", "协议动画 5 个 + 知识库机器人 MVP", "蒋献/诸葛斌"),
        ("9 月", "全部交付物初稿完成", "全体"),
        ("10 月", "中期总结材料汇总", "李传煌"),
        ("11-12 月", "中期工作总结 + 质量标准制定", "全体"),
    ]
    
    # 创建时间线
    y_start = 1.5
    for i, (month, task, owner) in enumerate(milestones):
        y = y_start + i * 0.7
        
        # 月份标签
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.5), Inches(y), Inches(1.2), Inches(0.5)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0x2E, 0x86, 0xC1)
        shape.line.fill.background()
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = month
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        p.alignment = PP_ALIGN.CENTER
        
        # 任务描述
        txBox = slide.shapes.add_textbox(Inches(2), Inches(y), Inches(5), Inches(0.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = task
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
        
        # 负责人
        txBox2 = slide.shapes.add_textbox(Inches(7.2), Inches(y), Inches(2), Inches(0.5))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = f"负责人：{owner}"
        p2.font.size = Pt(12)
        p2.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
        p2.alignment = PP_ALIGN.RIGHT
    
    return slide

def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # 1. 封面
    add_title_slide(
        prs, 
        "计算机网络\"101 计划\"课程建设",
        "第一次研讨会汇报"
    )
    
    # 2. 项目背景
    add_content_slide(prs, "一、项目背景", [
        "• 项目名称：计算机网络（模块一，本科，专业核心课）",
        "• 负责人：诸葛斌（浙江工商大学）",
        "• 团队成员：金蓉、高明、李传煌、蒋献",
        "• 建设周期：2026.01 - 2027.12（2 年）",
        "",
        "核心改革方向：",
        "• 以教育部\"101 计划\"为引领",
        "• 国产云平台为底座",
        "• 深度融合智能体（Agent）工具",
        "• 构建\"智能生成 - 私有答疑 - 智慧管理 - 实战部署\"全链条教学新范式",
    ])
    
    # 3. 官方进度安排
    add_content_slide(prs, "二、官方进度安排", [
        "2026 年：",
        "• 1-2 月：完成实施方案，组建专委会和工作组",
        "• 3-4 月：项目调研 + 研讨交流",
        "• 5-10 月：推进核心课程、教材、实践项目、师资团队重点任务建设",
        "• 11-12 月：中期工作总结，确立\"四个核心\"建设质量标准与认定办法",
        "",
        "2027 年底前：",
        "• 验收认定一批核心课程及配套教材、实践项目、师资团队",
        "• 建立年度评估及动态调整机制",
    ])
    
    # 4. 现有成果基础
    add_content_slide(prs, "三、现有成果基础", [
        "MOOC 在线课程：",
        "• 《高级网络通信原理实践》（中国大学 MOOC）",
        "• 228 人学习，第 14 周/共 18 周",
        "• 十三章结构（1-8 章园区网络，9-13 章网络虚拟化）",
        "• 与阿里云计算有限公司合作建立云实验室平台",
        "",
        "实验课程（省级线上一流课程）：",
        "• 《计算机网络实验》（浙江省在线开放课程共享平台）",
        "• 7 期/623 人次/22 所高校/59 万次访问",
        "• 负责人：金蓉",
        "",
        "已出版教材：",
        "• 《网络通信原理实践（微课视频版）》（清华大学出版社，2024.09）",
        "• 《系统级编程及分布式应用实现技术》（已定稿，预计 2026 年初出版）",
    ])
    
    # 5. 三大核心任务
    add_content_slide(prs, "四、三大核心任务", [
        "任务一：课程建设（负责人：诸葛斌）",
        "• 知识体系梳理（50-60 个关键知识点）",
        "• 知识点图谱构建",
        "• 课程建设指南编写",
        "• 教学大纲设计",
        "• 教学资源开发",
        "",
        "任务二：教材编写（负责人：诸葛斌 + 高明）",
        "• 出版社联络",
        "• 教材框架设计",
        "• 样章撰写",
        "• 案例资源开发",
        "• 教材撰写",
        "",
        "任务三：实践案例（负责人：诸葛斌 + 蒋献）",
        "• 行业企业联络",
        "• 实验项目设计",
        "• 智能体辅助实验",
        "• 实践案例群建设",
        "• 云端实验环境升级",
    ])
    
    # 6. 任务分工表
    add_content_slide(prs, "五、任务分工表", [
        [
            "序号", "子任务", "负责人", "配合人", "交付物", "时间节点"
        ],
        [
            "1.1", "知识体系梳理", "诸葛斌", "金蓉", "课程知识体系文档", "2026-06-30"
        ],
        [
            "1.2", "知识点图谱构建", "蒋献", "诸葛斌", "知识点图谱（可视化）", "2026-07-15"
        ],
        [
            "2.1", "教材框架设计", "高明", "金蓉", "教材框架设计文档", "2026-07-15"
        ],
        [
            "2.2", "样章撰写", "金蓉", "高明", "样章 1-2 章", "2026-07-31"
        ],
        [
            "3.1", "实验项目设计", "蒋献", "金蓉", "实验项目清单（≥10 个）", "2026-07-31"
        ],
        [
            "3.2", "实践案例群建设", "蒋献", "金蓉", "实践案例群（≥5 个）", "2026-09-30"
        ],
    ], is_table=True)
    
    # 7. 关键里程碑
    add_milestone_slide(prs)
    
    # 8. 协作机制
    add_content_slide(prs, "六、协作机制", [
        "会议制度：",
        "• 双周线上碰头会（周五 15:00-16:00）",
        "• 每季度线下研讨会（至少 2 次）",
        "",
        "文档管理：",
        "• 统一存放：钉钉群文件/云平台",
        "• 命名规范：YYYYMMDD_任务名称_负责人",
        "• 每月月底文档归档",
        "",
        "进度汇报：",
        "• 每周五 17:00 前提交本周进度报告",
        "• 遇到问题 24 小时内反馈",
    ])
    
    # 9. 6 月行动计划
    add_content_slide(prs, "七、6 月行动计划", [
        "第一周（6/1-6/7）：",
        "• ✅ 第一次研讨会",
        "• MOOC 十三章知识点初步梳理",
        "• 教材框架设计资料收集",
        "",
        "第二周（6/8-6/14）：",
        "• 知识点清单细化（50-60 个）",
        "• 教材框架设计初稿",
        "• 双周碰头会（6/12）",
        "",
        "第三周（6/15-6/21）：",
        "• 知识点图谱构建",
        "• 教学大纲修订",
        "• 协议动画选题（5 个优先协议）",
        "",
        "第四周（6/22-6/28）：",
        "• 知识体系文档整合",
        "• 教学大纲终稿",
        "• 双周碰头会（6/26）",
        "",
        "第五周（6/29-6/30）：",
        "• 6 月工作总结",
        "• 7 月计划制定",
    ])
    
    # 10. 会议决议
    add_content_slide(prs, "八、会议决议", [
        "1. 任务分工确认：",
        "   全体参会人员确认各自负责的任务模块，按《任务分工表》执行。",
        "",
        "2. 协作机制确认：",
        "   双周线上碰头会（周五 15:00-16:00），文档统一存放钉钉群文件/云平台。",
        "",
        "3. 时间节点确认：",
        "   按关键里程碑推进，每月月底提交进度报告。",
        "",
        "4. 下次会议：",
        "   2026 年 6 月 13 日（周五）15:00-16:00，第一次双周碰头会。",
    ])
    
    # 11. 结束页
    add_title_slide(
        prs,
        "感谢聆听",
        "敬请批评指正"
    )
    
    # 保存 PPT
    output_path = os.path.join(os.path.dirname(__file__), "计算机网络 101 计划_第一次研讨会汇报.pptx")
    prs.save(output_path)
    print(f"PPT 已生成：{output_path}")
    print(f"文件大小：{os.path.getsize(output_path) / 1024:.1f} KB")

if __name__ == "__main__":
    main()
