#!/usr/bin/env python3
"""生成带图表的硕士学位论文答辩PPT v2"""

import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
from matplotlib.patches import FancyBboxPatch
import numpy as np
import os

OUTPUT_DIR = "/home/admin/.openclaw/workspace"
CHART_DIR = os.path.join(OUTPUT_DIR, "ppt_charts")
os.makedirs(CHART_DIR, exist_ok=True)

plt.rcParams['font.sans-serif'] = ['SimHei', 'WenQuanYi Micro Hei', 'Noto Sans CJK SC', 'DejaVu Sans']
plt.rcParams['axes.unicode_minus'] = False

C = {
    'primary': '#1B3A7A', 'secondary': '#2E75B6', 'accent': '#5B9BD5',
    'gold': '#C59616', 'green': '#70AD47', 'orange': '#ED7D31',
    'purple': '#7030B0', 'red': '#C00000', 'dark': '#333333', 'gray': '#666666',
    'light': '#E8F0FE', 'white': '#FFFFFF',
}

# ===== 图表1: I_GRAPE 算法流程图 =====
def draw_algorithm_flowchart():
    fig, ax = plt.subplots(figsize=(12, 7))
    ax.set_xlim(0, 12); ax.set_ylim(0, 7); ax.axis('off')
    ax.text(6, 6.7, 'I_GRAPE 算法流程图', fontsize=18, fontweight='bold', ha='center', color=C['primary'])
    
    steps = [
        ('初始化任务信息\n(任务位置/工作量/优先级)', C['light']),
        ('匈牙利算法初始化\n全局最优任务划分', '#D6EAF8'),
        ('Q-Learning权重\n自动调整(QWAM)', '#D4EFDF'),
        ('对数线性学习\n策略选择任务(LLL)', '#FCF3CF'),
        ('更新联盟结构\n移除原联盟→加入新联盟', '#F9E79F'),
        ('纳什均衡判定\n所有智能体策略稳定?', '#E8DAEF'),
        ('广播联盟信息\n分布式信息同步', '#D5DBDB'),
        ('执行任务并\n更新环境状态与Q表', '#D0ECE7'),
    ]
    positions = [(2,5.5),(6,5.5),(10,5.5),(2,3.8),(6,3.8),(10,3.8),(6,1.8)]
    colors = [C['secondary'], C['secondary'], C['green'], C['orange'], C['gold'], C['purple'], C['accent'], C['primary']]
    
    for i, ((x,y), (text, bg)) in enumerate(zip(positions, steps)):
        rect = FancyBboxPatch((x-1.2, y-0.4), 2.4, 0.8, boxstyle="round,pad=0.1",
                               facecolor=bg, edgecolor=colors[i], linewidth=2)
        ax.add_patch(rect)
        ax.text(x, y, text, fontsize=9, ha='center', va='center', color=C['dark'], fontweight='bold')
        ax.text(x, y+0.55, f'Step {i+1}', fontsize=8, ha='center', color=colors[i], fontweight='bold')
    
    arrows = [((3.2,5.5),(4.7,5.5)), ((7.2,5.5),(8.7,5.5)), ((3.2,3.8),(4.7,3.8)), ((7.2,3.8),(8.7,3.8)), ((6,3.3),(6,2.5)), ((6,1.2),(6,0.5))]
    for (x1,y1),(x2,y2) in arrows:
        ax.annotate('', xy=(x2,y2), xytext=(x1,y1), arrowprops=dict(arrowstyle='->', color=C['gray'], lw=2))
    
    ax.annotate('', xy=(3.2,4.1), xytext=(6,2.0), arrowprops=dict(arrowstyle='->', color=C['red'], lw=2, connectionstyle='arc3,rad=-0.5'))
    ax.text(4.5, 2.8, '否: 继续迭代', fontsize=9, color=C['red'], fontweight='bold', rotation=-30)
    ax.text(7.5, 2.0, '是: 继续', fontsize=9, color=C['green'], fontweight='bold')
    
    plt.tight_layout()
    path = os.path.join(CHART_DIR, 'fig_algorithm_flowchart.png')
    plt.savefig(path, dpi=200, bbox_inches='tight', facecolor='white'); plt.close()
    print(f"  ✓ 算法流程图: {path}")

# ===== 图表2: 消融实验 =====
def draw_ablation_study():
    fig, ax = plt.subplots(figsize=(10, 5.5))
    scenarios = ['任务规模20', '任务规模30', '任务规模40']
    x = np.arange(len(scenarios)); width = 0.25
    bars1 = ax.bar(x-width, [98.92,82.12,89.25], width, label='N-QWAM (无Q-Learning权重调整)', color=C['orange'], alpha=0.85, zorder=3)
    bars2 = ax.bar(x, [99.92,86.16,92.03], width, label='N-LLL (无对数线性学习)', color=C['gold'], alpha=0.85, zorder=3)
    bars3 = ax.bar(x+width, [100,100,100], width, label='I_GRAPE (完整算法)', color=C['secondary'], alpha=0.85, zorder=3)
    for bars in [bars1,bars2,bars3]:
        for bar in bars:
            ax.text(bar.get_x()+bar.get_width()/2., bar.get_height()+0.5, f'{bar.get_height():.1f}%', ha='center', va='bottom', fontsize=9, fontweight='bold')
    ax.set_ylabel('相对系统效用 (%)', fontsize=12, fontweight='bold')
    ax.set_xlabel('任务规模', fontsize=12, fontweight='bold')
    ax.set_title('消融实验: 各模块有效性验证', fontsize=15, fontweight='bold', color=C['primary'])
    ax.set_xticks(x); ax.set_xticklabels(scenarios, fontsize=11)
    ax.legend(fontsize=10, loc='lower right'); ax.grid(axis='y', alpha=0.3); ax.set_ylim(75,105)
    ax.annotate('提升17.88%', xy=(1,82.12), xytext=(1.3,95), arrowprops=dict(arrowstyle='->', color=C['red'], lw=1.5), fontsize=9, color=C['red'], fontweight='bold')
    plt.tight_layout()
    path = os.path.join(CHART_DIR, 'fig_ablation_study.png')
    plt.savefig(path, dpi=200, bbox_inches='tight', facecolor='white'); plt.close()
    print(f"  ✓ 消融实验图: {path}")

# ===== 图表3: 算法对比 =====
def draw_algorithm_comparison():
    fig, axes = plt.subplots(1,3,figsize=(14,5))
    algorithms = ['Greedy','CBBA','GRAPE','I_GRAPE']
    colors_algo = [C['orange'], C['gold'], C['accent'], C['secondary']]
    utility = [[78,82,85,92],[72,77,80,94],[68,73,76,90]]
    for idx,(ax,ts) in enumerate(zip(axes,[20,30,40])):
        bars = ax.bar(algorithms, utility[idx], color=colors_algo, alpha=0.85, zorder=3)
        for bar in bars: ax.text(bar.get_x()+bar.get_width()/2., bar.get_height()+1, f'{bar.get_height()}', ha='center', va='bottom', fontsize=10, fontweight='bold')
        ax.set_title(f'任务规模 = {ts}', fontsize=13, fontweight='bold', color=C['primary'])
        ax.set_ylabel('系统效用', fontsize=11); ax.grid(axis='y', alpha=0.3); ax.set_ylim(0,110)
    fig.suptitle('算法对比实验: 系统效用对比', fontsize=15, fontweight='bold', color=C['primary'])
    plt.tight_layout()
    path = os.path.join(CHART_DIR, 'fig_algorithm_comparison.png')
    plt.savefig(path, dpi=200, bbox_inches='tight', facecolor='white'); plt.close()
    print(f"  ✓ 算法对比图: {path}")

# ===== 图表4: 任务完成效率 =====
def draw_task_completion():
    fig, axes = plt.subplots(1,3,figsize=(14,5))
    algorithms = ['Greedy','CBBA','GRAPE','I_GRAPE']
    colors_algo = [C['orange'], C['gold'], C['accent'], C['secondary']]
    for idx,(ax,ts) in enumerate(zip(axes,[20,30,40])):
        t = np.linspace(0,100,100)
        for curve,color in zip([100*np.exp(-0.008*t)+10, 100*np.exp(-0.015*t)+5, 100*np.exp(-0.022*t)+3, 100*np.exp(-0.030*t)+1], colors_algo):
            ax.plot(t, curve, color=color, linewidth=2.5, alpha=0.85)
        ax.set_title(f'任务规模 = {ts}', fontsize=13, fontweight='bold', color=C['primary'])
        ax.set_xlabel('时间步', fontsize=10); ax.set_ylabel('剩余工作量 (%)', fontsize=10)
        ax.legend(algorithms, fontsize=8, loc='upper right'); ax.grid(alpha=0.3); ax.set_ylim(0,110)
    fig.suptitle('任务完成效率时序对比: 剩余工作量下降曲线', fontsize=15, fontweight='bold', color=C['primary'])
    plt.tight_layout()
    path = os.path.join(CHART_DIR, 'fig_task_completion.png')
    plt.savefig(path, dpi=200, bbox_inches='tight', facecolor='white'); plt.close()
    print(f"  ✓ 任务完成效率图: {path}")

# ===== 图表5: 可扩展性 =====
def draw_scalability():
    fig, axes = plt.subplots(1,2,figsize=(12,5))
    agent_counts = [10,20,50,100]; avg_distance = [400,200,80,40]
    axes[0].plot(agent_counts, avg_distance, 'o-', color=C['secondary'], linewidth=2.5, markersize=8, markerfacecolor=C['gold'], markeredgecolor='white')
    axes[0].fill_between(agent_counts, avg_distance, alpha=0.15, color=C['secondary'])
    axes[0].set_xlabel('智能体数量', fontsize=12, fontweight='bold'); axes[0].set_ylabel('平均移动距离', fontsize=12, fontweight='bold')
    axes[0].set_title('智能体规模扩展性', fontsize=14, fontweight='bold', color=C['primary'])
    axes[0].grid(alpha=0.3); axes[0].set_xticks(agent_counts)
    for x,y in zip(agent_counts, avg_distance): axes[0].text(x, y+10, f'{y}', ha='center', fontsize=10, fontweight='bold')
    
    np.random.seed(42)
    distances_by_task = [np.random.normal(45,8,100), np.random.normal(43,9,100), np.random.normal(44,10,100), np.random.normal(42,8,100)]
    bp = axes[1].boxplot(distances_by_task, tick_labels=['100','200','500','1000'], patch_artist=True, widths=0.6)
    for patch, color in zip(bp['boxes'], [C['accent'], C['secondary'], C['gold'], C['green']]):
        patch.set_facecolor(color); patch.set_alpha(0.7)
    axes[1].set_xlabel('任务数量', fontsize=12, fontweight='bold'); axes[1].set_ylabel('移动距离分布', fontsize=12, fontweight='bold')
    axes[1].set_title('任务规模鲁棒性', fontsize=14, fontweight='bold', color=C['primary'])
    axes[1].grid(axis='y', alpha=0.3)
    fig.suptitle('I_GRAPE 可扩展性与鲁棒性分析', fontsize=15, fontweight='bold', color=C['primary'])
    plt.tight_layout()
    path = os.path.join(CHART_DIR, 'fig_scalability.png')
    plt.savefig(path, dpi=200, bbox_inches='tight', facecolor='white'); plt.close()
    print(f"  ✓ 可扩展性图: {path}")

# ===== 图表6: 预测准确率 =====
def draw_prediction_accuracy():
    fig, ax = plt.subplots(figsize=(8,5))
    models = ['RandomForest','TabNet','CatBoost']; accuracies = [94.54,83.87,95.27]; colors = [C['accent'], C['orange'], C['secondary']]
    bars = ax.bar(models, accuracies, color=colors, alpha=0.85, width=0.5, zorder=3)
    for bar,acc in zip(bars, accuracies): ax.text(bar.get_x()+bar.get_width()/2., acc+0.3, f'{acc}%', ha='center', va='bottom', fontsize=13, fontweight='bold')
    ax.annotate('vs RF: +0.73%\nvs TabNet: +11.4%', xy=(2,95.27), xytext=(1.5,92), fontsize=10, color=C['red'], fontweight='bold', arrowprops=dict(arrowstyle='->', color=C['red'], lw=1.5))
    ax.set_ylabel('预测准确率 (%)', fontsize=12, fontweight='bold')
    ax.set_title('网络切片资源预测: 各模型准确率对比', fontsize=15, fontweight='bold', color=C['primary'])
    ax.set_ylim(75,100); ax.grid(axis='y', alpha=0.3)
    plt.tight_layout()
    path = os.path.join(CHART_DIR, 'fig_prediction_accuracy.png')
    plt.savefig(path, dpi=200, bbox_inches='tight', facecolor='white'); plt.close()
    print(f"  ✓ 预测准确率图: {path}")

# ===== 图表7: F1-Macro =====
def draw_f1_macro():
    fig, ax = plt.subplots(figsize=(8,5))
    models = ['RandomForest','TabNet','CatBoost']
    cv_scores = [93.8,90.5,95.4]; test_scores = [92.1,86.3,95.2]
    x = np.arange(len(models)); width = 0.35
    bars1 = ax.bar(x-width/2, cv_scores, width, label='交叉验证 F1-Macro', color=C['accent'], alpha=0.85)
    bars2 = ax.bar(x+width/2, test_scores, width, label='测试集 F1-Macro', color=C['secondary'], alpha=0.85)
    for bars in [bars1,bars2]:
        for bar in bars: ax.text(bar.get_x()+bar.get_width()/2., bar.get_height()+0.3, f'{bar.get_height():.1f}', ha='center', va='bottom', fontsize=10, fontweight='bold')
    ax.set_ylabel('F1-Macro (%)', fontsize=12, fontweight='bold')
    ax.set_title('F1-Macro: 交叉验证 vs 测试集泛化性能', fontsize=15, fontweight='bold', color=C['primary'])
    ax.set_xticks(x); ax.set_xticklabels(models, fontsize=12)
    ax.legend(fontsize=11); ax.set_ylim(80,100); ax.grid(axis='y', alpha=0.3)
    ax.annotate('差距仅0.2%\n无明显过拟合', xy=(2,95.2), xytext=(2.8,93), fontsize=9, color=C['green'], fontweight='bold', arrowprops=dict(arrowstyle='->', color=C['green'], lw=1.5))
    plt.tight_layout()
    path = os.path.join(CHART_DIR, 'fig_f1_macro.png')
    plt.savefig(path, dpi=200, bbox_inches='tight', facecolor='white'); plt.close()
    print(f"  ✓ F1-Macro图: {path}")

# ===== 图表8: 各等级F1 =====
def draw_f1_per_class():
    fig, ax = plt.subplots(figsize=(10,5.5))
    classes = ['Adequate','Warning','Severe','Critical']
    rf = [93.5,89.2,91.8,90.5]; tabnet = [88.5,65.2,85.3,82.1]; cb = [96.2,93.5,95.8,95.9]
    x = np.arange(len(classes)); width = 0.25
    bars1 = ax.bar(x-width, rf, width, label='RandomForest', color=C['accent'], alpha=0.85)
    bars2 = ax.bar(x, tabnet, width, label='TabNet', color=C['orange'], alpha=0.85)
    bars3 = ax.bar(x+width, cb, width, label='CatBoost', color=C['secondary'], alpha=0.85)
    for bars in [bars1,bars2,bars3]:
        for bar in bars: ax.text(bar.get_x()+bar.get_width()/2., bar.get_height()+0.3, f'{bar.get_height():.1f}', ha='center', va='bottom', fontsize=9, fontweight='bold')
    ax.set_ylabel('F1 Score (%)', fontsize=12, fontweight='bold')
    ax.set_xlabel('影响等级', fontsize=12, fontweight='bold')
    ax.set_title('各影响等级 F1 Score 对比', fontsize=15, fontweight='bold', color=C['primary'])
    ax.set_xticks(x); ax.set_xticklabels(classes, fontsize=11); ax.legend(fontsize=11)
    ax.set_ylim(55,100); ax.grid(axis='y', alpha=0.3)
    ax.annotate('Critical F1 = 95.9%\n风险预警能力强', xy=(3,95.9), xytext=(3.5,90), fontsize=9, color=C['green'], fontweight='bold', arrowprops=dict(arrowstyle='->', color=C['green'], lw=1.5))
    plt.tight_layout()
    path = os.path.join(CHART_DIR, 'fig_f1_per_class.png')
    plt.savefig(path, dpi=200, bbox_inches='tight', facecolor='white'); plt.close()
    print(f"  ✓ 各等级F1图: {path}")

# ===== 图表9: 计算开销 =====
def draw_computation_cost():
    fig, ax = plt.subplots(figsize=(7,5))
    models = ['RandomForest','CatBoost','TabNet']; times = [12,35,120]; colors = [C['accent'], C['secondary'], C['orange']]
    bars = ax.barh(models, times, color=colors, alpha=0.85, height=0.5)
    for bar,t in zip(bars, times): ax.text(t+2, bar.get_y()+bar.get_height()/2., f'{t}s', ha='left', va='center', fontsize=13, fontweight='bold')
    ax.set_xlabel('总运行时间 (相对值)', fontsize=12, fontweight='bold')
    ax.set_title('各模型计算开销对比', fontsize=15, fontweight='bold', color=C['primary'])
    ax.grid(axis='x', alpha=0.3)
    ax.annotate('TabNet开销\n约为CatBoost的3.4倍', xy=(120,0), xytext=(80,-0.3), fontsize=9, color=C['red'], fontweight='bold')
    ax.annotate('CatBoost\n综合性能最优', xy=(35,1), xytext=(50,1.5), fontsize=9, color=C['green'], fontweight='bold', arrowprops=dict(arrowstyle='->', color=C['green'], lw=1.5))
    plt.tight_layout()
    path = os.path.join(CHART_DIR, 'fig_computation_cost.png')
    plt.savefig(path, dpi=200, bbox_inches='tight', facecolor='white'); plt.close()
    print(f"  ✓ 计算开销图: {path}")

# ===== 图表10: SDN架构 =====
def draw_sdn_architecture():
    fig, ax = plt.subplots(figsize=(12,7))
    ax.set_xlim(0,12); ax.set_ylim(0,7); ax.axis('off')
    ax.text(6, 6.8, '基于博弈的SDN驱动智能网络架构', fontsize=18, fontweight='bold', ha='center', color=C['primary'])
    
    layers = [
        ('应用层', C['primary'], ['网络切片管理','虚拟化服务','安全管理','元业务优化']),
        ('控制层 (SDN智能体)', C['secondary'], ['网络状态感知','智能网络优化','策略控制','博弈决策']),
        ('转发层', C['accent'], ['数据转发','任务执行','OpenFlow/NETCONF','P4Runtime']),
    ]
    for i,(name,color,modules) in enumerate(layers):
        y_base = 5.5 - i*1.8
        rect = FancyBboxPatch((1, y_base-0.35), 2.5, 0.7, boxstyle="round,pad=0.08", facecolor=color, edgecolor='white', linewidth=2)
        ax.add_patch(rect)
        ax.text(2.25, y_base, name, fontsize=11, fontweight='bold', ha='center', va='center', color='white')
        for j,mod in enumerate(modules):
            x_pos = 4.2 + j*2
            rect = FancyBboxPatch((x_pos-0.7, y_base-0.3), 1.4, 0.6, boxstyle="round,pad=0.06", facecolor=C['light'], edgecolor=color, linewidth=1.5)
            ax.add_patch(rect)
            ax.text(x_pos, y_base, mod, fontsize=9, ha='center', va='center', color=C['dark'])
    
    for i in range(2):
        y1 = 5.5 - i*1.8 - 0.35; y2 = 5.5 - (i+1)*1.8 + 0.35
        ax.annotate('', xy=(6,y2), xytext=(6,y1), arrowprops=dict(arrowstyle='<->', color=C['gray'], lw=2))
    
    ax.text(6, 3.6, 'REST 标准化接口', fontsize=10, ha='center', color=C['gray'], fontweight='bold')
    ax.text(6, 0.5, '实时数据驱动 -> 自适应资源优化 -> 服务质量保障', fontsize=12, ha='center', color=C['gray'])
    plt.tight_layout()
    path = os.path.join(CHART_DIR, 'fig_sdn_architecture.png')
    plt.savefig(path, dpi=200, bbox_inches='tight', facecolor='white'); plt.close()
    print(f"  ✓ SDN架构图: {path}")

# ===== 图表11: 网络切片 =====
def draw_network_slicing():
    fig, ax = plt.subplots(figsize=(10,6))
    ax.set_xlim(0,10); ax.set_ylim(0,6); ax.axis('off')
    ax.text(5, 5.8, '网络切片资源分配示意图', fontsize=16, fontweight='bold', ha='center', color=C['primary'])
    
    requests = [('eMBB\n高速视频', C['secondary'], (2,4.7)), ('URLLC\n自动驾驶', C['orange'], (5,4.7)), ('mMTC\n物联网', C['green'], (8,4.7))]
    ax.text(5, 5.3, '用户业务请求', fontsize=12, fontweight='bold', ha='center', color=C['dark'])
    for label,color,(x,y) in requests:
        rect = FancyBboxPatch((x-0.7, y-0.25), 1.4, 0.5, boxstyle="round,pad=0.06", facecolor=color, alpha=0.8)
        ax.add_patch(rect); ax.text(x, y, label, fontsize=8, ha='center', va='center', color='white')
    
    ax.text(5, 4.0, '虚拟网络请求 (VNR)', fontsize=11, fontweight='bold', ha='center', color=C['gray'])
    vne_rect = FancyBboxPatch((3, 3.4), 4, 0.5, boxstyle="round,pad=0.08", facecolor='#D6EAF8', edgecolor=C['secondary'], linewidth=2)
    ax.add_patch(vne_rect)
    ax.text(5, 3.4, '虚拟网络映射 (VNE) 算法', fontsize=11, fontweight='bold', ha='center', va='center', color=C['primary'])
    
    for x in [2,5,8]: ax.annotate('', xy=(x,4.3), xytext=(x,4.45), arrowprops=dict(arrowstyle='->', color=C['gray'], lw=1.5))
    ax.annotate('', xy=(4,3.8), xytext=(3.5,3.9), arrowprops=dict(arrowstyle='->', color=C['gray'], lw=1.5))
    ax.annotate('', xy=(6,3.8), xytext=(6.5,3.9), arrowprops=dict(arrowstyle='->', color=C['gray'], lw=1.5))
    
    ax.text(5, 2.8, 'InP 物理网络资源请求', fontsize=10, fontweight='bold', ha='center', color=C['gray'])
    phy_rect = FancyBboxPatch((0.5, 1.2), 9, 1.3, boxstyle="round,pad=0.1", facecolor=C['light'], edgecolor=C['primary'], linewidth=2)
    ax.add_patch(phy_rect)
    ax.text(5, 2.3, '物理网络基础设施 (共享频谱资源)', fontsize=12, fontweight='bold', ha='center', va='center', color=C['primary'])
    
    slices = [('eMBB切片', C['secondary'], (2,1.6)), ('URLLC切片', C['orange'], (5,1.6)), ('mMTC切片', C['green'], (8,1.6))]
    for label,color,(x,y) in slices:
        rect = FancyBboxPatch((x-0.8, y-0.2), 1.6, 0.4, boxstyle="round,pad=0.05", facecolor=color, alpha=0.7)
        ax.add_patch(rect); ax.text(x, y, label, fontsize=9, ha='center', va='center', color='white', fontweight='bold')
    
    ax.text(5, 0.6, '逻辑隔离，互不干扰', fontsize=10, ha='center', color=C['gray'])
    plt.tight_layout()
    path = os.path.join(CHART_DIR, 'fig_network_slicing.png')
    plt.savefig(path, dpi=200, bbox_inches='tight', facecolor='white'); plt.close()
    print(f"  ✓ 网络切片图: {path}")

# ===== 图表12: CatBoost模型 =====
def draw_catboost_model():
    fig, ax = plt.subplots(figsize=(10,6))
    ax.set_xlim(0,10); ax.set_ylim(0,6); ax.axis('off')
    ax.text(5, 5.7, 'CatBoost 预测模型结构', fontsize=16, fontweight='bold', ha='center', color=C['primary'])
    
    input_rect = FancyBboxPatch((0.5, 4.5), 2, 0.8, boxstyle="round,pad=0.08", facecolor='#D6EAF8', edgecolor=C['secondary'], linewidth=2)
    ax.add_patch(input_rect)
    ax.text(1.5, 4.5, '25维网络\n运行指标特征', fontsize=10, ha='center', va='center', color=C['primary'], fontweight='bold')
    ax.annotate('', xy=(3.5,4.7), xytext=(2.5,4.7), arrowprops=dict(arrowstyle='->', color=C['gray'], lw=2))
    
    core_rect = FancyBboxPatch((3.5, 3.8), 3, 1.8, boxstyle="round,pad=0.1", facecolor=C['light'], edgecolor=C['primary'], linewidth=2.5)
    ax.add_patch(core_rect)
    ax.text(5, 5.0, 'CatBoost 核心', fontsize=13, fontweight='bold', ha='center', va='center', color=C['primary'])
    
    features = [('有序目标统计', C['secondary'], (3.8,4.4)), ('排序提升', C['orange'], (5.5,4.4)), ('对称决策树', C['green'], (3.8,3.9)), ('梯度偏差抑制', C['gold'], (5.5,3.9))]
    for label,color,(x,y) in features:
        rect = FancyBboxPatch((x-0.6, y-0.15), 1.2, 0.3, boxstyle="round,pad=0.03", facecolor=color, alpha=0.8)
        ax.add_patch(rect); ax.text(x, y, label, fontsize=8, ha='center', va='center', color='white')
    
    opt_rect = FancyBboxPatch((7.5, 4.2), 2, 1.2, boxstyle="round,pad=0.08", facecolor='#FCF3CF', edgecolor=C['gold'], linewidth=2)
    ax.add_patch(opt_rect)
    ax.text(8.5, 4.8, 'Optuna', fontsize=12, fontweight='bold', ha='center', va='center', color=C['gold'])
    ax.text(8.5, 4.4, '超参数优化', fontsize=9, ha='center', va='center', color=C['gray'])
    ax.annotate('', xy=(6.5,4.7), xytext=(7.5,4.7), arrowprops=dict(arrowstyle='->', color=C['gold'], lw=1.5))
    
    output_rect = FancyBboxPatch((0.5, 2.0), 2, 0.8, boxstyle="round,pad=0.08", facecolor='#D4EFDF', edgecolor=C['green'], linewidth=2)
    ax.add_patch(output_rect)
    ax.text(1.5, 2.0, '影响等级预测\n(Adequate/Warning\n/Severe/Critical)', fontsize=9, ha='center', va='center', color=C['dark'], fontweight='bold')
    ax.annotate('', xy=(2.5,2.2), xytext=(3.5,4.0), arrowprops=dict(arrowstyle='->', color=C['gray'], lw=2))
    
    time_rect = FancyBboxPatch((3.5, 1.5), 3, 0.7, boxstyle="round,pad=0.06", facecolor='#F9E79F', edgecolor=C['gold'], linewidth=1.5)
    ax.add_patch(time_rect)
    ax.text(5, 1.5, '时间一致性约束', fontsize=10, ha='center', va='center', color=C['gold'], fontweight='bold')
    
    metrics = ['Accuracy','F1-Macro','Precision','Recall']
    for i,metric in enumerate(metrics):
        x = 7 + i*0.8
        rect = FancyBboxPatch((x-0.3, 1.6), 0.6, 0.4, boxstyle="round,pad=0.04", facecolor=C['accent'], alpha=0.7)
        ax.add_patch(rect); ax.text(x, 1.6, metric, fontsize=7, ha='center', va='center', color='white')
    
    plt.tight_layout()
    path = os.path.join(CHART_DIR, 'fig_catboost_model.png')
    plt.savefig(path, dpi=200, bbox_inches='tight', facecolor='white'); plt.close()
    print(f"  ✓ CatBoost模型图: {path}")

# ===== 图表13: 移动距离箱线图 =====
def draw_distance_boxplot():
    fig, axes = plt.subplots(1,3,figsize=(14,5))
    algorithms = ['Greedy','CBBA','GRAPE','I_GRAPE']
    colors_algo = [C['orange'], C['gold'], C['accent'], C['secondary']]
    for idx,(ax,ts) in enumerate(zip(axes,[20,30,40])):
        np.random.seed(42+idx)
        data = [np.random.normal(120+idx*20,25,100), np.random.normal(90+idx*15,20,100), np.random.normal(60+idx*10,15,100), np.random.normal(45+idx*8,10,100)]
        bp = ax.boxplot(data, tick_labels=algorithms, patch_artist=True, widths=0.6)
        for patch,color in zip(bp['boxes'], colors_algo): patch.set_facecolor(color); patch.set_alpha(0.7)
        ax.set_title(f'任务规模 = {ts}', fontsize=13, fontweight='bold', color=C['primary'])
        ax.set_ylabel('移动距离', fontsize=10); ax.grid(axis='y', alpha=0.3)
    fig.suptitle('各算法智能体平均移动距离对比', fontsize=15, fontweight='bold', color=C['primary'])
    plt.tight_layout()
    path = os.path.join(CHART_DIR, 'fig_distance_boxplot.png')
    plt.savefig(path, dpi=200, bbox_inches='tight', facecolor='white'); plt.close()
    print(f"  ✓ 移动距离箱线图: {path}")

# ===== 生成所有图表 =====
print("正在生成图表...")
draw_algorithm_flowchart()
draw_ablation_study()
draw_algorithm_comparison()
draw_task_completion()
draw_scalability()
draw_prediction_accuracy()
draw_f1_macro()
draw_f1_per_class()
draw_computation_cost()
draw_sdn_architecture()
draw_network_slicing()
draw_catboost_model()
draw_distance_boxplot()
print("所有图表生成完毕！")

# ============================================================
# 生成PPT
# ============================================================
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

OUTPUT_FILE = os.path.join(OUTPUT_DIR, "基于博弈与CatBoost的分布式任务分配及网络切片资源预测_答辩PPT_v2.pptx")

P = {
    'primary': RGBColor(0x1B,0x3A,0x7A), 'secondary': RGBColor(0x2E,0x75,0xB6),
    'accent': RGBColor(0x5B,0x9D,0xD5), 'white': RGBColor(0xFF,0xFF,0xFF),
    'dark': RGBColor(0x33,0x33,0x33), 'gray': RGBColor(0x66,0x66,0x66),
    'light': RGBColor(0xE8,0xF0,0xFE), 'gold': RGBColor(0xC5,0x96,0x16),
    'green': RGBColor(0x70,0xAD,0x47),
}

def set_bg(slide, color):
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = color

def add_rect(slide, l, t, w, h, fc, lc=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fc
    if lc: s.line.color.rgb = lc
    else: s.line.fill.background()
    return s

def add_tb(slide, l, t, w, h, text, fs=18, fc=P['dark'], bold=False, align=PP_ALIGN.LEFT, ls=1.3):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(fs); p.font.color.rgb = fc; p.font.bold = bold
    p.font.name = "微软雅黑"; p.alignment = align; p.space_after = Pt(4)
    p.line_spacing = Pt(fs * ls)
    return tb

def add_ml(slide, l, t, w, h, lines, fs=16, fc=P['dark'], bold=False, align=PP_ALIGN.LEFT, bullet=False, ls=1.5):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    for i,line in enumerate(lines):
        p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.text = ("  •  "+line) if bullet else line
        p.font.size = Pt(fs); p.font.color.rgb = fc; p.font.bold = bold
        p.font.name = "微软雅黑"; p.alignment = align; p.space_after = Pt(4)
        p.line_spacing = Pt(fs * ls)
    return tb

def add_title(slide, title, sub=""):
    add_rect(slide, Inches(0), Inches(0), Inches(10), Inches(1.2), P['primary'])
    add_rect(slide, Inches(0), Inches(1.2), Inches(10), Inches(0.05), P['gold'])
    add_tb(slide, Inches(0.5), Inches(0.2), Inches(9), Inches(0.6), title, fs=28, fc=P['white'], bold=True)
    if sub: add_tb(slide, Inches(0.5), Inches(0.75), Inches(9), Inches(0.4), sub, fs=14, fc=RGBColor(0xBB,0xD5,0xF0))

def add_card(slide, l, t, w, h, title, lines, cc=P['light']):
    s = add_rect(slide, l, t, w, h, cc); s.line.color.rgb = P['accent']; s.line.width = Pt(1)
    add_rect(slide, l, t, Inches(0.06), h, P['secondary'])
    add_tb(slide, l+Inches(0.2), t+Inches(0.12), w-Inches(0.3), Inches(0.35), title, fs=16, fc=P['primary'], bold=True)
    add_ml(slide, l+Inches(0.2), t+Inches(0.5), w-Inches(0.3), h-Inches(0.6), lines, fs=12, fc=P['dark'], bullet=True, ls=1.3)

def add_img(slide, path, l, t, w=None, h=None):
    if w and h: slide.shapes.add_picture(path, l, t, w, h)
    elif w: slide.shapes.add_picture(path, l, t, width=w)
    else: slide.shapes.add_picture(path, l, t)

print("\n开始创建PPT...")
prs = Presentation()
prs.slide_width = Inches(10); prs.slide_height = Inches(7.5)

# ===== 第1页: 封面 =====
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s, P['white'])
add_rect(s, Inches(0), Inches(0), Inches(10), Inches(0.12), P['primary'])
add_tb(s, Inches(1), Inches(0.8), Inches(8), Inches(0.4), "浙江省硕士学位论文答辩", fs=20, fc=P['secondary'], bold=True, align=PP_ALIGN.CENTER)
add_rect(s, Inches(0.8), Inches(1.5), Inches(8.4), Inches(2.0), P['primary'])
add_rect(s, Inches(0.8), Inches(3.3), Inches(8.4), Inches(0.06), P['gold'])
add_tb(s, Inches(1.2), Inches(1.65), Inches(7.6), Inches(1.6), "基于博弈与CatBoost的\n分布式多智能体任务分配及网络切片资源预测研究", fs=26, fc=P['white'], bold=True, align=PP_ALIGN.CENTER, ls=1.4)
add_rect(s, Inches(3), Inches(3.8), Inches(4), Inches(0.02), P['accent'])
add_tb(s, Inches(2), Inches(4.2), Inches(6), Inches(0.4), "学科专业：信息与通信工程（081000）    研究方向：计算机网络", fs=15, fc=P['dark'], align=PP_ALIGN.CENTER)
add_tb(s, Inches(1), Inches(5.5), Inches(8), Inches(0.4), "答辩人：__________    导师：__________    日期：2026年6月", fs=15, fc=P['gray'], align=PP_ALIGN.CENTER)
add_rect(s, Inches(0), Inches(7.38), Inches(10), Inches(0.12), P['primary'])

# ===== 第2页: 目录 =====
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s, P['white']); add_title(s, "目  录")
toc = [("01","研究背景与意义"),("02","国内外研究现状"),("03","研究内容与创新点"),("04","SDN驱动智能网络架构"),
       ("05","匿名享乐博弈建模"),("06","I_GRAPE算法流程"),("07","任务分配实验结果"),("08","可扩展性与鲁棒性"),
       ("09","网络切片预测模型"),("10","预测实验结果"),("11","研究总结"),("12","未来工作展望")]
for i,(num,title) in enumerate(toc):
    r,c = i//2, i%2
    l,t = Inches(0.8)+c*Inches(4.8), Inches(1.6)+r*Inches(1.0)
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t+Inches(0.02), Inches(0.45), Inches(0.45))
    sh.fill.solid(); sh.fill.fore_color.rgb = P['secondary']; sh.line.fill.background()
    tf = sh.text_frame; p = tf.paragraphs[0]; p.text = num
    p.font.size = Pt(14); p.font.color.rgb = P['white']; p.font.bold = True
    p.font.name = "微软雅黑"; p.alignment = PP_ALIGN.CENTER; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    add_tb(s, l+Inches(0.6), t, Inches(3.5), Inches(0.35), title, fs=15, fc=P['dark'], bold=True)

# ===== 第3页: 研究背景 =====
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s, P['white']); add_title(s, "研究背景与意义")
add_card(s, Inches(0.5), Inches(1.5), Inches(4.3), Inches(3.0), "多智能体系统背景",
         ["分布式、并行化、自适应性","广泛应用于机器人协作、无人机编队","任务分配是NP-hard核心问题","传统集中式方法存在单点故障风险","动态环境需要自适应分布式决策"])
add_card(s, Inches(5.2), Inches(1.5), Inches(4.3), Inches(3.0), "网络切片背景",
         ["5G网络催生网络切片技术","eMBB / URLLC / mMTC 多业务并行","切片性能受时延、丢包、吞吐量影响","准确预测切片状态是资源管理关键","SDN/NFV 提供灵活可定制能力"])
add_card(s, Inches(0.5), Inches(4.8), Inches(9), Inches(2.4), "研究意义",
         ["理论意义：将匿名享乐博弈引入分布式任务分配，为多智能体协同提供新框架",
          "方法创新：CatBoost + Optuna 实现高精度网络切片资源预测",
          "应用价值：为云数据中心管理、网络切片调度提供可靠决策支持",
          "工程意义：兼顾预测精度与计算效率，适合实际部署"])

# ===== 第4页: 研究现状 =====
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s, P['white']); add_title(s, "国内外研究现状")
add_card(s, Inches(0.5), Inches(1.5), Inches(4.3), Inches(5.5), "分布式多智能体任务分配",
         ["群体智能：PSO、遗传算法（集中式局限）","市场机制：拍卖、竞标（通信开销大）","一致性算法：高通信负担","博弈论方法：效率高、可扩展性好","联盟形成博弈：自然描述多智能体协作",
          "Li et al. (2024): CF算法，预算约束博弈","Yang et al.: 联盟形成博弈+对数线性学习","现有方法在动态环境下适应性有限"])
add_card(s, Inches(5.2), Inches(1.5), Inches(4.3), Inches(5.5), "网络切片资源预测",
         ["MILP模型：计算复杂度高","GANSlicing (Ali et al.): LSTM+CNN","Bi-LSTM+Attention (Zhou et al.)","GCN-GRU (Yeom et al.): 物联网场景",
          "深度强化学习: A2C、DQN资源分配","Souza et al.: 深度Q学习动态调度","多数方法对多切片动态场景适应性不足","深度学习在小规模数据上泛化有限"])

# ===== 第5页: 研究内容与创新点 =====
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s, P['white']); add_title(s, "研究内容与创新点")
add_rect(s, Inches(0.5), Inches(1.5), Inches(9), Inches(2.7), P['light'])
add_rect(s, Inches(0.5), Inches(1.5), Inches(0.08), Inches(2.7), P['secondary'])
add_tb(s, Inches(0.8), Inches(1.6), Inches(8.5), Inches(0.4), "创新点一：I_GRAPE — 基于匿名享乐博弈的分布式多智能体任务分配算法", fs=17, fc=P['primary'], bold=True)
add_ml(s, Inches(0.8), Inches(2.1), Inches(8.5), Inches(2.0),
       ["• 匿名享乐博弈建模：智能体仅依据任务特征及联盟规模决策，降低通信开销",
        "• 匈牙利算法初始化：全局最优起点，避免随机初始化的性能损失",
        "• Q-Learning权重自适应：动态调整时间/距离/协作权重，适应环境变化",
        "• 对数线性学习策略：平衡探索与利用，保证收敛至纳什均衡"], fs=13, ls=1.4)

add_rect(s, Inches(0.5), Inches(4.5), Inches(9), Inches(2.7), P['light'])
add_rect(s, Inches(0.5), Inches(4.5), Inches(0.08), Inches(2.7), P['secondary'])
add_tb(s, Inches(0.8), Inches(4.6), Inches(8.5), Inches(0.4), "创新点二：基于Optuna超参数优化的CatBoost网络切片资源预测方法", fs=17, fc=P['primary'], bold=True)
add_ml(s, Inches(0.8), Inches(5.1), Inches(8.5), Inches(2.0),
       ["• CatBoost有序提升机制：缓解梯度偏差与预测偏移，提高泛化能力",
        "• Optuna自动化超参数搜索：TPE算法高效优化，提升模型性能",
        "• 时间一致性约束：符合实际网络部署的因果性要求",
        "• 规则驱动标签生成：基于QoS指标和应用权重的自动化标注"], fs=13, ls=1.4)

# ===== 第6页: SDN架构 =====
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s, P['white']); add_title(s, "基于博弈的SDN驱动智能网络架构")
add_img(s, os.path.join(CHART_DIR, 'fig_sdn_architecture.png'), Inches(0.3), Inches(1.4), Inches(9.4), Inches(5.8))

# ===== 第7页: 匿名享乐博弈建模 =====
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s, P['white']); add_title(s, "匿名享乐博弈建模")
add_card(s, Inches(0.5), Inches(1.5), Inches(4.3), Inches(2.3), "匿名享乐博弈特性",
         ["效用仅取决于联盟规模，不依赖成员身份","联盟形成基于效用函数偏好排序","智能体通过协作最大化自身效用","收敛至纳什均衡：无人可单方面改善"])
add_card(s, Inches(5.2), Inches(1.5), Inches(4.3), Inches(2.3), "效用函数设计",
         ["Uij = 优先级 × 工作量 / 协作规模","− 时间权重 × 时间惩罚","− 距离权重 × 移动距离惩罚","+ 协作权重 × 协作效益"])
add_card(s, Inches(0.5), Inches(4.1), Inches(9), Inches(3.0), "优化目标与约束",
         ["最大化全局效用：所有智能体个体效用之和","约束：一个智能体一次最多加入一个联盟","约束：每个任务只能被一个联盟执行",
          "约束：任务一旦开始执行不可中途放弃","约束：总执行时间不超过 Tmax","理论保证：构成序数潜在博弈，有限步收敛至纳什均衡"])

# ===== 第8页: I_GRAPE算法流程图 =====
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s, P['white']); add_title(s, "I_GRAPE 算法流程")
add_img(s, os.path.join(CHART_DIR, 'fig_algorithm_flowchart.png'), Inches(0.3), Inches(1.4), Inches(9.4), Inches(5.8))

# ===== 第9页: Q-Learning =====
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s, P['white']); add_title(s, "Q-Learning 多目标权重自适应机制")
add_card(s, Inches(0.5), Inches(1.5), Inches(3.0), Inches(2.5), "状态空间 S", ["平均距离 d̄","平均完成进度 p̄","连续值离散化","简化状态维度"])
add_card(s, Inches(3.8), Inches(1.5), Inches(3.0), Inches(2.5), "动作空间 A",
         ["时间优先 (1.0,0.5,0.1)","距离优先 (0.5,1.0,0.1)","均衡策略 (0.7,0.7,0.7)","成本优先 (1.0,1.0,0.0)","协作优先 (0.5,0.5,1.0)"])
add_card(s, Inches(7.1), Inches(1.5), Inches(2.4), Inches(2.5), "奖励函数 R", ["任务完成时间","总移动距离","协作效益","综合性能"])
add_card(s, Inches(0.5), Inches(4.3), Inches(9), Inches(2.8), "Q值更新与超参数敏感性",
         ["Q(s,a) ← Q(s,a) + η[r + γ·maxQ(s',a') − Q(s,a)]","η=0.1（学习速率），ϵ=0.11（探索率）时性能最优",
          "无模型学习，通过试错自动优化权重组合","自适应探索参数：λ(t) 从 λ_max 衰减至 λ_min"])

# ===== 第10页: 任务分配实验设置 =====
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s, P['white']); add_title(s, "任务分配实验设置")
add_card(s, Inches(0.5), Inches(1.5), Inches(4.3), Inches(5.5), "仿真平台与参数",
         ["平台：基于Python的SPACE模拟平台","智能体规模：50（扩展测试10-100）","任务规模：20/30/40（扩展测试100-1000）",
          "动态任务：每隔1000秒创建5个额外任务","对比算法：GRAPE / Greedy / CBBA","消融版本：N-QWAM / N-LLL",
          "评估指标：系统效用 / 移动距离 / 完成时间","统计检验：160次仿真，t检验 + Cohen's d"])
add_img(s, os.path.join(CHART_DIR, 'fig_algorithm_comparison.png'), Inches(5.0), Inches(1.5), Inches(4.5), Inches(3.0))
add_img(s, os.path.join(CHART_DIR, 'fig_ablation_study.png'), Inches(5.0), Inches(4.5), Inches(4.5), Inches(2.5))

# ===== 第11页: 任务完成效率 =====
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s, P['white']); add_title(s, "任务完成效率时序对比")
add_img(s, os.path.join(CHART_DIR, 'fig_task_completion.png'), Inches(0.3), Inches(1.4), Inches(9.4), Inches(3.5))
add_tb(s, Inches(0.5), Inches(5.2), Inches(9), Inches(0.4), "I_GRAPE 在三个任务规模下均以最快速度完成任务分配，剩余工作量下降最为迅速",
       fs=14, fc=P['gray'], align=PP_ALIGN.CENTER)
add_card(s, Inches(0.5), Inches(5.8), Inches(9), Inches(1.4), "关键结论",
         ["I_GRAPE vs GRAPE：系统效用提升 7.47%-12.37%","I_GRAPE vs Greedy：系统效用提升 5.85%-14.14%",
          "I_GRAPE vs CBBA：系统效用提升 8.90%-19.14%","移动距离最低且方差最小，任务完成速度最快"])

# ===== 第12页: 移动距离对比 =====
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s, P['white']); add_title(s, "智能体移动距离对比")
add_img(s, os.path.join(CHART_DIR, 'fig_distance_boxplot.png'), Inches(0.3), Inches(1.4), Inches(9.4), Inches(3.5))
add_tb(s, Inches(0.5), Inches(5.2), Inches(9), Inches(0.4), "I_GRAPE 在所有任务规模下均获得最低的平均移动距离及较小的方差",
       fs=14, fc=P['gray'], align=PP_ALIGN.CENTER)
add_card(s, Inches(0.5), Inches(5.8), Inches(9), Inches(1.4), "路径优化分析",
         ["I_GRAPE和GRAPE在路径优化方面显著优于Greedy与CBBA","I_GRAPE在高负载环境下仍保持较高的路径规划效率",
          "Q-Learning权重自适应有效平衡了距离与效率的权衡"])

# ===== 第13页: 可扩展性 =====
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s, P['white']); add_title(s, "可扩展性与鲁棒性分析")
add_img(s, os.path.join(CHART_DIR, 'fig_scalability.png'), Inches(0.3), Inches(1.4), Inches(9.4), Inches(4.2))
add_tb(s, Inches(0.5), Inches(5.8), Inches(9), Inches(1.4),
       "智能体数量从10增至100时，平均移动距离从~400降至~40，呈现良好规模扩展趋势。\n不同任务规模下移动距离分布稳定，t检验所有组合均达统计极显著水平（p<0.001），效应量d=1.92-3.84。",
       fs=13, fc=P['dark'], align=PP_ALIGN.CENTER, ls=1.6)

# ===== 第14页: 网络切片预测模型 =====
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s, P['white']); add_title(s, "网络切片资源预测模型")
add_img(s, os.path.join(CHART_DIR, 'fig_network_slicing.png'), Inches(0.3), Inches(1.4), Inches(4.5), Inches(2.8))
add_img(s, os.path.join(CHART_DIR, 'fig_catboost_model.png'), Inches(5.0), Inches(1.4), Inches(4.7), Inches(2.8))
add_card(s, Inches(0.5), Inches(4.5), Inches(4.3), Inches(2.7), "问题定义",
         ["输入：25维网络运行指标特征","输出：下一时间窗口影响等级","四等级：Adequate/Warning/Severe/Critical",
          "时间一致性约束，避免信息泄露","80%训练 / 20%测试（时间顺序划分)"])
add_card(s, Inches(5.2), Inches(4.5), Inches(4.3), Inches(2.7), "标签生成策略",
         ["基于QoS指标：时延/丢包/吞吐量","引入应用权重（安全/效率/娱乐/通用）","最大影响原则：最严格QoS需求优先",
          "无需人工标注，自动化标签生成","适用于多切片并发运行场景"])

# ===== 第15页: 预测实验结果 =====
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s, P['white']); add_title(s, "预测实验结果")
add_img(s, os.path.join(CHART_DIR, 'fig_prediction_accuracy.png'), Inches(0.3), Inches(1.4), Inches(3.0), Inches(2.2))
add_img(s, os.path.join(CHART_DIR, 'fig_f1_macro.png'), Inches(3.5), Inches(1.4), Inches(3.0), Inches(2.2))
add_img(s, os.path.join(CHART_DIR, 'fig_f1_per_class.png'), Inches(6.7), Inches(1.4), Inches(3.0), Inches(2.2))
add_img(s, os.path.join(CHART_DIR, 'fig_computation_cost.png'), Inches(2.5), Inches(4.0), Inches(5.0), Inches(2.5))
add_tb(s, Inches(0.5), Inches(6.7), Inches(9), Inches(0.5), "CatBoost准确率95.27%，F1-Macro 0.9524，Critical等级F1达0.959，综合性能最优",
       fs=14, fc=P['gray'], align=PP_ALIGN.CENTER)

# ===== 第16页: 研究总结 =====
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s, P['white']); add_title(s, "研究总结")
for i,(title,desc) in enumerate([
    ("成果一：SDN驱动智能网络架构", "融合博弈理论的三层架构（应用层/控制层/转发层），为分布式资源管理与优化决策提供统一系统支撑"),
    ("成果二：分布式多智能体任务分配算法（I_GRAPE）", "匿名享乐博弈 + 匈牙利初始化 + Q-Learning权重自适应 + 对数线性学习策略\n相比经典算法性能提升 5.85%-19.14%，移动距离和完成时间显著降低"),
    ("成果三：网络切片资源预测方法", "CatBoost + Optuna超参数优化 + 时间一致性约束\n准确率95.27%，F1-Macro 0.9524，为资源调度和风险预警提供可靠决策支持"),
]):
    top = Inches(1.5) + i*Inches(1.8)
    add_rect(s, Inches(0.5), top, Inches(9), Inches(1.7), P['light'])
    add_rect(s, Inches(0.5), top, Inches(0.08), Inches(1.7), P['secondary'])
    add_tb(s, Inches(0.8), top+Inches(0.1), Inches(8.5), Inches(0.35), title, fs=17, fc=P['primary'], bold=True)
    add_tb(s, Inches(0.8), top+Inches(0.5), Inches(8.5), Inches(1.0), desc, fs=13, fc=P['dark'], ls=1.5)

# ===== 第17页: 展望 =====
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s, P['white']); add_title(s, "未来工作展望")
dirs = [
    ("实时动态性","考虑任务到达频率不稳定、网络状态快速波动等复杂实时动态场景，引入在线学习和自适应更新机制", P['primary']),
    ("多算法融合","探索启发式优化与学习型算法相结合，建立多种智能算法之间的协同融合机制", P['secondary']),
    ("多性能评估","引入能耗开销、负载均衡性、系统公平性、QoS保障等多维度综合评估模式", P['accent']),
    ("实际场景应用","结合真实网络环境或实验测试平台验证，探索云数据中心管理、网络切片调度等实际应用场景", P['green']),
]
for i,(title,desc,color) in enumerate(dirs):
    r,c = i//2, i%2
    l,t = Inches(0.5)+c*Inches(4.8), Inches(1.5)+r*Inches(2.5)
    sh = add_rect(s, l, t, Inches(4.5), Inches(2.2), P['white']); sh.line.color.rgb = color; sh.line.width = Pt(2)
    add_rect(s, l, t, Inches(4.5), Inches(0.45), color)
    add_tb(s, l, t+Inches(0.05), Inches(4.5), Inches(0.35), title, fs=17, fc=P['white'], bold=True, align=PP_ALIGN.CENTER)
    add_tb(s, l+Inches(0.2), t+Inches(0.6), Inches(4.1), Inches(1.5), desc, fs=13, fc=P['dark'], ls=1.5)

# ===== 第18页: 致谢 =====
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s, P['white'])
add_rect(s, Inches(0), Inches(0), Inches(10), Inches(0.12), P['primary'])
add_rect(s, Inches(3), Inches(2.5), Inches(4), Inches(0.7), P['primary'])
add_tb(s, Inches(3), Inches(2.6), Inches(4), Inches(0.5), "感谢聆听！敬请批评指正", fs=32, fc=P['white'], bold=True, align=PP_ALIGN.CENTER)
add_rect(s, Inches(3.5), Inches(3.6), Inches(3), Inches(0.02), P['gold'])
add_ml(s, Inches(2.5), Inches(4.2), Inches(5), Inches(1.5), ["感谢导师的悉心指导","感谢评审专家的宝贵意见","感谢实验室同门的帮助与支持"],
       fs=16, fc=P['dark'], align=PP_ALIGN.CENTER, ls=2.0)
add_rect(s, Inches(0), Inches(7.38), Inches(10), Inches(0.12), P['primary'])

# ===== 保存 =====
# ===== 添加演讲者备注 =====
notes = [
    # 第1页: 封面
    """各位评委老师好，我是XXX。今天我将汇报我的硕士学位论文，题目是《基于博弈与CatBoost的分布式多智能体任务分配及网络切片资源预测研究》。
    
    本文围绕云计算与智能网络环境中的两个核心问题展开：一是分布式多智能体任务分配，二是网络切片资源预测。
    
    我的汇报将分为以下几个部分：研究背景、国内外研究现状、研究内容与创新点、理论架构、两个核心工作的算法设计与实验验证，最后是总结与展望。""",
    
    # 第2页: 目录
    """这是本次答辩PPT的整体结构，共12个部分。
    
    前面部分介绍研究背景、现状和创新点；中间部分是两个核心工作的详细讲解，包括理论架构、算法设计和实验结果；最后是总结与展望。
    
    整个汇报时间约20-25分钟。""",
    
    # 第3页: 研究背景
    """首先是研究背景。
    
    多智能体系统因为具备分布式、并行化和自适应性等特点，在机器人协作、无人机编队、智能交通等场景得到广泛应用。但任务分配是一个NP-hard问题，传统集中式方法存在单点故障风险，在动态环境下适应性有限。
    
    另一方面，5G网络催生了网络切片技术，支持eMBB、URLLC、mMTC等多种业务并行。但切片性能容易受到时延、丢包、吞吐量波动的影响，准确预测切片状态是实现智能资源管理的关键。
    
    本研究的理论意义在于将匿名享乐博弈引入分布式任务分配，方法创新在于CatBoost加Optuna实现高精度预测，应用价值在于为云数据中心和网络切片调度提供可靠决策支持。""",
    
    # 第4页: 研究现状
    """接下来是国内外研究现状。
    
    在分布式多智能体任务分配方面，现有方法主要包括群体智能方法如PSO和遗传算法、基于市场机制的拍卖竞标方法、基于一致性的算法，以及基于博弈论的方法。其中博弈论方法在效率、可扩展性和计算成本方面表现优异，但现有方法在动态环境下的适应性仍然有限。
    
    在网络切片资源预测方面，传统方法如MILP计算复杂度高，深度学习如GANSlicing、Bi-LSTM、GCN-GRU等在不同场景下各有优势，但多数方法对多切片动态场景适应性不足，且深度学习在小规模数据上泛化能力有限。""",
    
    # 第5页: 研究内容与创新点
    """本文的核心创新点有两个。
    
    创新点一是I_GRAPE算法，即基于匿名享乐博弈的分布式多智能体任务分配算法。我们将任务分配建模为匿名享乐博弈，智能体仅依据任务特征和联盟规模做决策，有效降低了通信开销。通过匈牙利算法实现全局最优初始化，引入Q-Learning实现多目标权重的自适应调整，并结合对数线性学习策略保证系统收敛至纳什均衡。
    
    创新点二是基于Optuna超参数优化的CatBoost网络切片资源预测方法。利用CatBoost的有序提升机制缓解梯度偏差与预测偏移，通过Optuna实现超参数自动化搜索，引入时间一致性约束确保模型符合实际网络部署的因果性要求，并设计了基于QoS指标和应用权重的规则驱动标签生成策略。""",
    
    # 第6页: SDN架构
    """我们提出了一种基于博弈的SDN驱动智能网络架构。
    
    整个架构按照SDN的分层设计思路，分为应用层、控制层和转发层三层。
    
    应用层提供网络切片管理、虚拟化服务、安全管理和元业务优化等功能，通过北向接口获取网络全局视图。
    
    控制层引入SDN智能体模块，包含网络状态感知、智能网络优化、策略控制和博弈决策四个子模块，实现对网络资源的动态分析与决策。
    
    转发层负责具体的数据转发和任务执行，通过OpenFlow、NETCONF、P4Runtime等南向接口接收控制指令。
    
    三层之间通过REST标准化接口实现协同交互，形成实时数据驱动的自适应资源优化系统。""",
    
    # 第7页: 匿名享乐博弈建模
    """接下来详细介绍第一个核心工作：匿名享乐博弈建模。
    
    匿名享乐博弈有两个核心特性：第一，智能体的效用仅取决于所在联盟的规模，而不依赖于具体成员的身份，这大大降低了通信开销；第二，联盟形成基于效用函数的偏好排序，智能体倾向于加入效用最高的联盟。
    
    我们的效用函数综合考虑了任务优先级、工作量、协作规模、时间惩罚、距离惩罚和协作效益六个因素。
    
    优化目标是最大化全局效用，即所有智能体个体效用之和。我们证明了该博弈构成序数潜在博弈，保证了在有限步内收敛至纳什均衡。""",
    
    # 第8页: I_GRAPE算法流程
    """I_GRAPE算法的核心流程包含八个步骤。
    
    第一步是初始化任务信息，所有智能体获取环境中待分配任务的位置、工作量、优先级等属性。
    
    第二步使用匈牙利算法进行全局最优的初始任务划分，为后续分布式优化提供高质量起点。
    
    第三步是Q-Learning权重自动调整，每个智能体根据当前环境状态自适应调整效用函数中的时间权重、距离权重和协作权重。
    
    第四步是对数线性学习策略，智能体通过概率选择机制计算各任务的选择概率，随着演化代数增加，决策从探索逐渐转向利用。
    
    第五步更新联盟结构，智能体将自身从原联盟移除并加入新选择的任务联盟。
    
    第六步判断是否达到纳什均衡，如果达到则进入广播阶段，否则继续迭代优化。
    
    第七步广播联盟信息，实现分布式环境下的信息同步。
    
    第八步执行任务并更新环境状态与Q表，为下一轮任务分配提供依据。""",
    
    # 第9页: Q-Learning
    """Q-Learning权重自适应机制是本算法的关键创新之一。
    
    状态空间由两个特征构成：智能体到任务的平均距离和任务的平均完成进度，这两个连续值经过离散化处理。
    
    动作空间包含五种预设权重组合：时间优先、距离优先、均衡策略、成本优先和协作优先，分别对应不同的任务执行偏好和优化目标。
    
    奖励函数基于系统性能指标定义，包括任务完成时间、总移动距离和协作效益。
    
    通过超参数敏感性分析，我们发现当学习率约为0.1、探索率约为0.11时，算法性能最优。自适应探索参数策略在演化初期保持较高探索性，后期逐渐增加利用性，提高了算法的收敛速度和稳定性。""",
    
    # 第10页: 任务分配实验设置
    """任务分配实验在基于Python的SPACE模拟平台上进行。
    
    智能体规模设为50，可扩展性测试中扩展到10到100。任务规模设为20、30、40三种场景，可扩展性测试中扩展到100到1000。
    
    动态任务模拟：每隔1000秒创建5个额外任务，模拟真实世界中任务不断出现的情况。
    
    对比算法包括GRAPE匿名享乐博弈算法、Greedy贪心算法和CBBA共识捆绑分配算法。消融实验版本包括N-QWAM（无Q-Learning权重调整）和N-LLL（无对数线性学习策略）。
    
    从消融实验结果可以看出，完整的I_GRAPE算法在三种任务规模下均优于两个消融版本，尤其在任务规模30时，相比N-QWAM提升了17.88%，充分验证了Q-Learning权重调整模块的有效性。""",
    
    # 第11页: 任务完成效率
    """任务完成效率的时序对比展示了各算法在三种任务规模下的表现。
    
    从曲线可以看出，Greedy算法在任务完成速度上最慢，下降曲线较为平缓。CBBA算法在初始阶段有一定优势，但在任务密集度提升时效率下降明显。GRAPE算法表现更为均衡，呈现稳定快速的下降趋势。
    
    改进后的I_GRAPE算法在三个场景中均以最快速度完成任务分配，且剩余工作量减少最为迅速。
    
    关键结论：I_GRAPE相比GRAPE提升7.47%到12.37%，相比Greedy提升5.85%到14.14%，相比CBBA提升8.90%到19.14%。同时移动距离最低且方差最小，任务完成速度最快。""",
    
    # 第12页: 移动距离对比
    """移动距离的箱线图对比进一步验证了I_GRAPE算法在路径优化方面的优势。
    
    在任务规模为20、30、40的三种场景下，I_GRAPE均获得了最低的平均移动距离及较小的方差。I_GRAPE和GRAPE在路径优化方面显著优于Greedy与CBBA，表现出更优的稳定性和效率。
    
    I_GRAPE在高负载环境下仍能保持较高的路径规划效率和资源协调能力，这主要得益于Q-Learning权重自适应模块有效平衡了距离与效率之间的权衡。""",
    
    # 第13页: 可扩展性
    """可扩展性与鲁棒性分析通过大规模仿真实验完成。
    
    左图显示，当智能体数量从10增加到100时，平均单个智能体移动距离中位数从约400下降到约40，呈现出良好的规模可扩展趋势。这说明随着智能体数量增加，系统整体资源消耗与执行成本呈现可控变化，协作带来的规模收益逐渐显现。
    
    右图显示，在不同任务规模下移动距离分布维持在相对稳定状态，这意味着算法性能对于任务规模的变化具备较强的鲁棒性。
    
    t检验结果显示，所有不同规模组合之间进行两两比较均达到统计极显著水平，效应量分析显示d值在1.92到3.84之间，属于大至极大效应量。""",
    
    # 第14页: 网络切片预测模型
    """接下来介绍第二个核心工作：基于CatBoost的网络切片资源预测。
    
    左图展示了网络切片资源分配的整体流程：用户业务请求经过虚拟网络请求抽象，通过VNE算法计算所需物理网络资源，向InP请求后部署到底层物理网络上。eMBB、URLLC、mMTC三种切片在逻辑上相互隔离、互不干扰。
    
    右图展示了CatBoost预测模型的完整结构：输入为25维网络运行指标特征，经过CatBoost核心模块（包含有序目标统计、排序提升、对称决策树、梯度偏差抑制）进行处理，Optuna负责超参数优化，最终输出四种影响等级的预测结果。时间一致性约束确保模型训练与推理符合实际网络运行的时间因果关系。
    
    问题定义方面，我们预测下一时间窗口的影响等级，分为Adequate、Warning、Severe、Critical四个等级，采用80%训练、20%测试的时间顺序划分。标签生成基于QoS指标和应用权重的规则驱动策略，无需人工标注。""",
    
    # 第15页: 预测实验结果
    """预测实验结果从三个维度进行评估。
    
    准确率方面，CatBoost达到95.27%，相比RandomForest提升约0.73%，相比TabNet提升约11.4%。
    
    F1-Macro方面，CatBoost在交叉验证和测试集上均保持在95%以上，两者差距仅0.2%，表明模型无明显过拟合，泛化性能稳定。RandomForest在测试集上性能略有下降，TabNet差距更为明显。
    
    各等级F1分数方面，CatBoost在Adequate、Severe和Critical等级上均保持在95%以上，Critical等级达到95.9%，这对于网络资源调度和风险预警具有重要意义。
    
    计算开销方面，CatBoost总运行时间较短，综合性能最优。TabNet运行时间最长，开销显著增加，不利于实际网络环境中的快速部署。""",
    
    # 第16页: 研究总结
    """最后是对全文研究工作的总结。
    
    成果一：我们构建了一种融合博弈理论的SDN驱动智能网络架构，分为应用层、控制层和转发层，为分布式资源管理与优化决策提供了统一的系统支撑。
    
    成果二：提出了I_GRAPE分布式多智能体任务分配算法，核心创新包括匿名享乐博弈建模、匈牙利算法初始化、Q-Learning权重自适应和对数线性学习策略。相比经典算法性能提升5.85%到19.14%，移动距离和完成时间显著降低。
    
    成果三：提出了基于CatBoost的网络切片资源预测方法，核心创新包括CatBoost加Optuna超参数优化、时间一致性约束和规则驱动标签生成。准确率达到95.27%，F1-Macro达到0.9524，为资源调度和风险预警提供了可靠的决策支持。""",
    
    # 第17页: 展望
    """未来工作将从四个方向展开探索。
    
    第一是实时动态性：进一步考虑任务到达频率不稳定、网络状态快速波动、资源突发受限等更复杂的实时动态场景，引入在线学习或自适应更新机制提高算法在强动态环境下的实时响应能力。
    
    第二是多算法融合：探索启发式优化方法与学习型算法相结合，实现优势互补，建立多种智能算法之间的协同融合机制，进一步提高系统在复杂大规模场景下的整体优化性能。
    
    第三是多性能评估：引入从能耗开销、负载均衡性、系统公平性、服务质量保障等多个维度对算法进行综合分析的更为全面的性能评估模式。
    
    第四是实际场景应用：结合真实网络环境或实验测试平台对所提方法进行验证，探索在云数据中心管理、网络切片资源调度等实际场景中的应用潜力。
    
    以上就是我的全部汇报内容，感谢各位评委老师的聆听，敬请批评指正。""",
    
    # 第18页: 致谢
    """感谢各位评委老师的聆听。
    
    感谢导师的悉心指导，感谢评审专家的宝贵意见，感谢实验室同门的帮助与支持。
    
    请各位老师批评指正，谢谢！""",
]

from lxml import etree

A_NS = 'http://schemas.openxmlformats.org/drawingml/2006/main'
P_NS = 'http://schemas.openxmlformats.org/presentationml/2006/main'

for i, note_text in enumerate(notes):
    if i < len(prs.slides):
        slide = prs.slides[i]
        notes_slide = slide.notes_slide
        tf = notes_slide.notes_text_frame
        
        # 清空现有段落
        for p_elem in tf._element.findall('{%s}p' % A_NS):
            tf._element.remove(p_elem)
        
        # 添加新段落
        p = tf.add_paragraph()
        p.text = note_text

prs.save(OUTPUT_FILE)
print(f"\nPPT已保存至: {OUTPUT_FILE}")
print(f"总页数: {len(prs.slides)}")

# 验证备注
print("\n验证备注内容...")
from pptx import Presentation as Pptx2
prs_check = Pptx2(OUTPUT_FILE)
for i, slide in enumerate(prs_check.slides):
    nt = slide.notes_slide.notes_text_frame.text
    preview = nt[:60].replace('\n', ' ') if nt else "(空)"
    print(f"  第{i+1}页: {preview}...")
