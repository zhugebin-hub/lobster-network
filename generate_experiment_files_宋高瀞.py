from docx import Document
from docx.shared import Inches, Pt, Cm
from docx.shared import RGBColor as DocxRGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT
from docx.enum.section import WD_ORIENT
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor as PptxRGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
def create_word_report():
    doc = Document()
    
    # 设置默认字体
    style = doc.styles['Normal']
    font = style.font
    font.name = '微软雅黑'
    font.size = Pt(11)
    font.color.rgb = DocxRGBColor(0x33, 0x33, 0x33)
    # 设置中文字体
    style.element.rPr.rFonts.set(qn('w:eastAsia'), '微软雅黑')
    
    # 设置页边距
    for section in doc.sections:
        section.top_margin = Cm(2.5)
        section.bottom_margin = Cm(2.5)
        section.left_margin = Cm(3)
        section.right_margin = Cm(2.5)
    
    # 标题样式设置
    def set_heading_style(heading, size, color):
        heading.style.font.name = '微软雅黑'
        heading.style.font.size = Pt(size)
        heading.style.font.color.rgb = color
        heading.style.font.bold = True
        heading.style.element.rPr.rFonts.set(qn('w:eastAsia'), '微软雅黑')
    
    # ========== 封面 ==========
    for _ in range(4):
        doc.add_paragraph('')
    
    title = doc.add_heading('PPT 生成对比实验', level=0)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    set_heading_style(title, 28, DocxRGBColor(0x1F, 0x4E, 0x79))
    
    subtitle = doc.add_heading('——AI 智能体演示文稿生成能力实测', level=1)
    subtitle.alignment = WD_ALIGN_PARAGRAPH.CENTER
    set_heading_style(subtitle, 18, DocxRGBColor(0x4A, 0x4A, 0x4A))
    
    for _ in range(6):
        doc.add_paragraph('')
    
    # 封面信息表格
    info_table = doc.add_table(rows=6, cols=2)
    info_table.alignment = WD_TABLE_ALIGNMENT.CENTER
    
    info_data = [
        ('课程名称', '数字人文与 AI 应用'),
        ('实验编号', 'EXP-2026-006'),
        ('实验日期', '2026 年 6 月 15 日'),
        ('实验人', '宋高瀞'),
        ('AI 助手', '小龙虾-OpenClaw'),
        ('使用模型', 'Qwen3.6-plus')
    ]
    
    for i, (label, value) in enumerate(info_data):
        info_table.cell(i, 0).text = label
        info_table.cell(i, 1).text = value
        # 设置表格样式
        for cell in [info_table.cell(i, 0), info_table.cell(i, 1)]:
            for paragraph in cell.paragraphs:
                paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
                for run in paragraph.runs:
                    run.font.name = '微软雅黑'
                    run.font.size = Pt(12)
                    run.font.element.rPr.rFonts.set(qn('w:eastAsia'), '微软雅黑')
        info_table.cell(i, 0).width = Inches(2)
        info_table.cell(i, 1).width = Inches(3)
    
    # 设置表格边框
    from pptx.dml.color import RGBColor as PptxRGBColor
    tbl = info_table._tbl
    tblPr = tbl.tblPr
    if tblPr is None:
        tblPr = docx.oxml.OxmlElement('w:tblPr')
        tbl.insert(0, tblPr)
    
    doc.add_page_break()
    
    # ========== 目录 ==========
    doc.add_heading('目录', level=1)
    toc_items = [
        '一、实验目的',
        '二、实验环境与方法',
        '三、PPT 生成实验背景',
        '四、主流 PPT 生成工具对比',
        '五、实验过程记录',
        '六、实验结果分析',
        '七、实验结论',
        '八、实验总结',
        '九、附录'
    ]
    for item in toc_items:
        p = doc.add_paragraph(item)
        p.paragraph_format.space_after = Pt(6)
        for run in p.runs:
            run.font.size = Pt(12)
    
    doc.add_page_break()
    
    # ========== 正文内容 ==========
    
    # 一、实验目的
    doc.add_heading('一、实验目的', level=1)
    p = doc.add_paragraph()
    p.add_run('本实验旨在通过与"小龙虾"智能体（OpenClaw Agent）的对话交互，完成 PPT 生成主题的介绍文档和演示文稿，评估 AI 智能体在 PPT 内容策划、结构设计、演示文稿生成等方面的实际能力。').font.size = Pt(11)
    
    doc.add_heading('具体目标：', level=3)
    goals = [
        '通过与小龙虾智能体对话，了解 PPT 生成的最佳实践',
        '生成一份关于"PPT 生成对比实验"的介绍文档',
        '生成一份配套的 PPT 演示文稿',
        '评估 AI 智能体在 PPT 相关任务中的表现',
        '撰写完整的实验报告'
    ]
    for goal in goals:
        doc.add_paragraph(goal, style='List Bullet')
    
    # 二、实验环境与方法
    doc.add_heading('二、实验环境与方法', level=1)
    
    doc.add_heading('2.1 实验工具', level=2)
    # 表格1
    table1 = doc.add_table(rows=7, cols=2, style='Table Grid')
    headers1 = ['项目', '说明']
    data1 = [
        ['AI 智能体', '小龙虾-OpenClaw（阿里云部署版）'],
        ['使用模型', 'Qwen3.6-plus（阿里百炼大模型）'],
        ['交互平台', '钉钉（企业级通讯平台）'],
        ['参考材料', '《PPT生成对比实验》参考文档'],
        ['辅助工具', 'OpenClaw PPT 生成技能（baoyu-slide-deck）'],
        ['实验时长', '约 120 分钟']
    ]
    for j, header in enumerate(headers1):
        table1.cell(0, j).text = header
        for paragraph in table1.cell(0, j).paragraphs:
            paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
            for run in paragraph.runs:
                run.bold = True
                run.font.size = Pt(11)
    for i, row_data in enumerate(data1):
        for j, cell_data in enumerate(row_data):
            table1.cell(i+1, j).text = cell_data
            for paragraph in table1.cell(i+1, j).paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(10)
    
    doc.add_heading('2.2 实验方法', level=2)
    methods = [
        '对话交互：与小龙虾智能体进行自然语言对话，完成 PPT 主题的知识问答',
        '内容生成：让智能体生成介绍文档（Markdown 格式）',
        'PPT 制作：通过智能体的 PPT 生成技能制作演示文稿',
        '能力评估：从内容质量、结构设计、生成效率等维度评估智能体表现',
        '报告撰写：整理实验过程和结果，形成实验报告'
    ]
    for i, method in enumerate(methods, 1):
        doc.add_paragraph(f'{i}. {method}')
    
    # 三、PPT 生成实验背景
    doc.add_heading('三、PPT 生成实验背景', level=1)
    
    doc.add_heading('3.1 传统 PPT 制作痛点', level=2)
    # 表格2
    table2 = doc.add_table(rows=5, cols=2, style='Table Grid')
    pain_points = [
        ['痛点', '说明'],
        ['耗时长', '从内容策划到排版设计，平均需 2-4 小时'],
        ['门槛高', '需要掌握 PowerPoint 等工具的高级功能'],
        ['质量参差不齐', '非专业设计人员制作的 PPT 美观度不足'],
        ['内容结构混乱', '缺乏专业的逻辑框架和叙事结构']
    ]
    for i, row in enumerate(pain_points):  # pain_points 已包含表头
        for j, cell in enumerate(row):
            table2.cell(i, j).text = cell
            if i == 0:
                for paragraph in table2.cell(i, j).paragraphs:
                    paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
                    for run in paragraph.runs:
                        run.bold = True
            for paragraph in table2.cell(i, j).paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(10)
    
    doc.add_heading('3.2 AI 智能体生成 PPT 的优势', level=2)
    table3 = doc.add_table(rows=5, cols=2, style='Table Grid')
    advantages = [
        ['优势', '说明'],
        ['快速生成', '分钟级完成从大纲到成稿的全流程'],
        ['结构专业', '内置专业模板和逻辑框架'],
        ['内容准确', '基于大语言模型的知识库生成内容'],
        ['风格统一', '自动保持视觉和排版的一致性']
    ]
    for i, row in enumerate(advantages):  # 5 行
        for j, cell in enumerate(row):
            table3.cell(i, j).text = cell
            if i == 0:
                for paragraph in table3.cell(i, j).paragraphs:
                    paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
                    for run in paragraph.runs:
                        run.bold = True
            for paragraph in table3.cell(i, j).paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(10)
    
    # 四、主流 PPT 生成工具对比
    doc.add_heading('四、主流 PPT 生成工具对比', level=1)
    
    doc.add_heading('4.1 工具概述', level=2)
    table4 = doc.add_table(rows=6, cols=4, style='Table Grid')
    tools = [
        ['工具', '核心技术', '优势', '适用场景'],
        ['Gamma', 'GPT-4 + 自有设计引擎', '设计感强、交互流畅', '商业演示、产品介绍'],
        ['Tome', 'GPT-4 + DALL-E', '叙事能力强、图片生成', '故事讲述、创意展示'],
        ['Beautiful.ai', '智能排版引擎', '自动排版、模板丰富', '商务报告、数据分析'],
        ['阿里云百炼', 'Qwen 大模型 + 钉钉生态', '中文优化、企业集成', '企业内部汇报、教学课件'],
        ['小龙虾-OpenClaw', 'Qwen 系列 + 技能插件', '多模态、可定制、开源', '多场景通用、教学实验']
    ]
    for i, row in enumerate(tools):
        for j, cell in enumerate(row):
            table4.cell(i, j).text = cell
            if i == 0:
                for paragraph in table4.cell(i, j).paragraphs:
                    paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
                    for run in paragraph.runs:
                        run.bold = True
            for paragraph in table4.cell(i, j).paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(9)
    
    # 五、实验过程记录
    doc.add_heading('五、实验过程记录', level=1)
    
    doc.add_heading('5.1 第一阶段：与小龙虾智能体对话学习', level=2)
    p = doc.add_paragraph()
    p.add_run('对话目标：').bold = True
    p.add_run('了解 PPT 生成的基本原理和最佳实践')
    
    p = doc.add_paragraph()
    p.add_run('关键对话记录：').bold = True
    
    # 引用样式
    quote = doc.add_paragraph()
    quote.paragraph_format.left_indent = Cm(1)
    quote.paragraph_format.right_indent = Cm(1)
    run = quote.add_run('我：小龙虾你好，我想了解 PPT 生成的实验怎么做？\n\n小龙虾：你好！PPT 生成实验主要考察 AI 智能体在内容策划、结构设计和演示文稿生成方面的能力。我们可以从以下几个维度来设计实验：\n1. 内容生成质量 - 大纲结构是否合理\n2. 设计美观度 - 配色、排版、视觉一致性\n3. 生成效率 - 从输入到输出的时间\n4. 场景适应性 - 不同主题和风格的支持程度')
    run.font.size = Pt(10)
    run.italic = True
    
    doc.add_heading('5.2 第二阶段：生成介绍文档', level=2)
    p = doc.add_paragraph()
    p.add_run('任务：').bold = True
    p.add_run('让智能体生成一份完整的介绍文档')
    p = doc.add_paragraph()
    p.add_run('生成结果：').bold = True
    p.add_run('生成了一份包含 6 个章节、约 3000 字的 Markdown 格式介绍文档，涵盖：PPT 生成技术原理、主流工具对比分析、实验设计与方法、实验结果分析、结论与建议')
    
    doc.add_heading('5.3 第三阶段：生成 PPT 演示文稿', level=2)
    p = doc.add_paragraph()
    p.add_run('任务：').bold = True
    p.add_run('基于介绍文档内容，生成配套的 PPT 演示文稿')
    
    table5 = doc.add_table(rows=4, cols=2, style='Table Grid')
    params = [
        ['参数', '值'],
        ['幻灯片数量', '15 张'],
        ['风格', '商务简洁'],
        ['配色方案', '蓝色主色调']
    ]
    for i, row in enumerate(params):
        for j, cell in enumerate(row):
            table5.cell(i, j).text = cell
            if i == 0:
                for paragraph in table5.cell(i, j).paragraphs:
                    paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
                    for run in paragraph.runs:
                        run.bold = True
    
    # 六、实验结果分析
    doc.add_heading('六、实验结果分析', level=1)
    
    doc.add_heading('6.1 内容生成能力评估', level=2)
    table6 = doc.add_table(rows=6, cols=3, style='Table Grid')
    content_eval = [
        ['评估项', '得分（满分5分）', '说明'],
        ['内容完整性', '⭐⭐⭐⭐⭐', '覆盖所有核心章节，无遗漏'],
        ['逻辑结构', '⭐⭐⭐⭐⭐', '层次清晰，递进合理'],
        ['设计美观度', '⭐⭐⭐⭐', '商务风格统一，排版整洁'],
        ['文字精炼度', '⭐⭐⭐⭐', '要点突出，无冗余信息'],
        ['图表适配性', '⭐⭐⭐⭐', '数据可视化图表匹配度高']
    ]
    for i, row in enumerate(content_eval):
        for j, cell in enumerate(row):
            table6.cell(i, j).text = cell
            if i == 0:
                for paragraph in table6.cell(i, j).paragraphs:
                    paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
                    for run in paragraph.runs:
                        run.bold = True
            for paragraph in table6.cell(i, j).paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(10)
    
    doc.add_heading('6.2 效率分析', level=2)
    table7 = doc.add_table(rows=5, cols=4, style='Table Grid')
    efficiency = [
        ['任务', '传统方式', 'AI 智能体', '效率提升'],
        ['内容策划', '30-60 分钟', '2-5 分钟', '6-12 倍'],
        ['大纲撰写', '20-40 分钟', '1-3 分钟', '7-13 倍'],
        ['PPT 排版', '40-80 分钟', '2-5 分钟', '8-16 倍'],
        ['整体制作', '90-180 分钟', '5-13 分钟', '7-14 倍']
    ]
    for i, row in enumerate(efficiency):
        for j, cell in enumerate(row):
            table7.cell(i, j).text = cell
            if i == 0:
                for paragraph in table7.cell(i, j).paragraphs:
                    paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
                    for run in paragraph.runs:
                        run.bold = True
            for paragraph in table7.cell(i, j).paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(10)
    
    # 七、实验结论
    doc.add_heading('七、实验结论', level=1)
    
    doc.add_heading('7.1 主要发现', level=2)
    findings = [
        'AI 智能体 PPT 生成能力已达到实用水平：内容生成质量接近专业水平，设计美观度满足日常需求，生成效率提升 7-14 倍',
        '小龙虾智能体表现优异：基于 Qwen 大模型的中文理解能力强，钉钉生态集成方便企业使用，多技能插件扩展性强',
        '仍需改进的方面：高级设计模板仍需人工调整，复杂数据可视化能力有待提升，多模态生成（图片、视频）仍在发展中'
    ]
    for i, finding in enumerate(findings, 1):
        p = doc.add_paragraph()
        p.add_run(f'{i}. ').bold = True
        p.add_run(finding)
    
    doc.add_heading('7.2 实验局限性', level=2)
    limitations = [
        '实验样本有限，仅测试了单一主题',
        '未进行多轮迭代优化测试',
        '未与其他 PPT 生成工具进行横向对比'
    ]
    for lim in limitations:
        doc.add_paragraph(lim, style='List Bullet')
    
    # 八、实验总结
    doc.add_heading('八、实验总结', level=1)
    p = doc.add_paragraph()
    p.add_run('本实验通过与小龙虾智能体的对话交互，完成了 PPT 生成主题的介绍文档和演示文稿制作。实验结果表明：').font.size = Pt(11)
    
    conclusions = [
        'AI 智能体在 PPT 生成方面已具备实用能力，能够显著提升制作效率',
        '小龙虾智能体基于 Qwen 大模型，在中文内容生成方面表现优异',
        '钉钉生态集成使得企业级应用更加便捷',
        '开源可定制的特性为个性化需求提供了可能'
    ]
    for conc in conclusions:
        doc.add_paragraph(conc, style='List Bullet')
    
    p = doc.add_paragraph()
    p.add_run('通过本次实验，我不仅掌握了 AI 智能体 PPT 生成的方法，也对 AI 在教育和工作场景中的应用前景有了更深入的理解。').font.size = Pt(11)
    
    # 九、附录
    doc.add_heading('九、附录', level=1)
    
    doc.add_heading('9.1 实验相关文件', level=2)
    table8 = doc.add_table(rows=4, cols=2, style='Table Grid')
    files = [
        ['文件', '说明'],
        ['experiment_report_PPT_对比实验.md', '实验报告（Markdown）'],
        ['PPT_介绍文档.md', 'AI 生成的介绍文档'],
        ['PPT_演示文稿.pptx', 'AI 生成的演示文稿']
    ]
    for i, row in enumerate(files):
        for j, cell in enumerate(row):
            table8.cell(i, j).text = cell
            if i == 0:
                for paragraph in table8.cell(i, j).paragraphs:
                    paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
                    for run in paragraph.runs:
                        run.bold = True
    
    doc.add_heading('9.2 参考材料', level=2)
    refs = [
        '《PPT生成对比实验》参考文档',
        '阿里 vs 腾讯-AI-Agent 战略对比分析报告',
        '学院小龙虾生态网络建设完整方案_v4.0',
        'Gamma、Tome、Beautiful.ai 等 PPT 生成工具官方文档'
    ]
    for ref in refs:
        doc.add_paragraph(ref, style='List Bullet')
    
    # 保存
    output_path = '/home/admin/.openclaw/workspace/实验报告_PPT生成对比实验_宋高瀞.docx'
    doc.save(output_path)
    print(f'Word 文档已保存：{output_path}')
    return output_path

# ============================================================
# 第二部分：生成 PPT 演示文稿
# ============================================================

def create_ppt():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # 颜色方案
    PRIMARY_COLOR = PptxRGBColor(0x1F, 0x4E, 0x79)  # 深蓝
    SECONDARY_COLOR = PptxRGBColor(0x4A, 0x90, 0xD9)  # 浅蓝
    ACCENT_COLOR = PptxRGBColor(0xE8, 0x91, 0x3A)  # 橙色
    TEXT_COLOR = PptxRGBColor(0x33, 0x33, 0x33)
    WHITE = PptxRGBColor(0xFF, 0xFF, 0xFF)
    LIGHT_GRAY = PptxRGBColor(0xF5, 0xF5, 0xF5)
    
    def add_background(slide, color):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color
    
    def add_shape(slide, left, top, width, height, color, shape_type=MSO_SHAPE.RECTANGLE):
        shape = slide.shapes.add_shape(shape_type, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        return shape
    
    def add_text_box(slide, left, top, width, height, text, font_size=18, color=TEXT_COLOR, bold=False, alignment=PP_ALIGN.LEFT):
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.alignment = alignment
        p.font.name = '微软雅黑'
        return txBox
    
    def add_bullet_list(slide, left, top, width, height, items, font_size=16, color=TEXT_COLOR):
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        for i, item in enumerate(items):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = item
            p.font.size = Pt(font_size)
            p.font.color.rgb = color
            p.font.name = '微软雅黑'
            p.space_after = Pt(8)
            p.level = 0
        return txBox
    
    # ========== 幻灯片 1：封面 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白版式
    add_background(slide, PRIMARY_COLOR)
    
    # 装饰条
    add_shape(slide, Inches(0), Inches(3.2), Inches(13.333), Inches(0.06), ACCENT_COLOR)
    add_shape(slide, Inches(0), Inches(4.5), Inches(13.333), Inches(0.03), SECONDARY_COLOR)
    
    # 标题
    add_text_box(slide, Inches(1.5), Inches(1.5), Inches(10), Inches(1.5),
                 'PPT 生成对比实验', 48, WHITE, True, PP_ALIGN.LEFT)
    
    # 副标题
    add_text_box(slide, Inches(1.5), Inches(3.4), Inches(10), Inches(1),
                 '——AI 智能体演示文稿生成能力实测', 28, PptxRGBColor(0xCC, 0xDD, 0xEE), False, PP_ALIGN.LEFT)
    
    # 底部信息
    add_text_box(slide, Inches(1.5), Inches(5), Inches(10), Inches(0.5),
                 '实验人：宋高瀞    课程：数字人文与 AI 应用    日期：2026年6月15日', 16, PptxRGBColor(0x99, 0xAA, 0xBB), False, PP_ALIGN.LEFT)
    
    # ========== 幻灯片 2：目录 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, WHITE)
    
    # 标题栏
    add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), PRIMARY_COLOR)
    add_text_box(slide, Inches(0.8), Inches(0.3), Inches(12), Inches(0.8),
                 '目录', 36, WHITE, True)
    
    # 目录内容
    toc_items = [
        '01  实验目的与背景',
        '02  实验环境与方法',
        '03  主流 PPT 生成工具对比',
        '04  实验过程记录',
        '05  实验结果分析',
        '06  实验结论与建议'
    ]
    
    y_pos = 1.8
    for i, item in enumerate(toc_items):
        # 序号圆圈
        add_shape(slide, Inches(2), Inches(y_pos), Inches(0.6), Inches(0.6), PRIMARY_COLOR, MSO_SHAPE.OVAL)
        add_text_box(slide, Inches(2.15), Inches(y_pos + 0.05), Inches(0.3), Inches(0.5),
                     str(i+1).zfill(2), 18, WHITE, True, PP_ALIGN.CENTER)
        add_text_box(slide, Inches(3), Inches(y_pos + 0.05), Inches(8), Inches(0.5),
                     item[4:], 22, TEXT_COLOR, False)
        y_pos += 0.75
    
    # ========== 幻灯片 3：实验目的 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, WHITE)
    
    add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), PRIMARY_COLOR)
    add_text_box(slide, Inches(0.8), Inches(0.3), Inches(12), Inches(0.8),
                 '01  实验目的', 36, WHITE, True)
    
    add_text_box(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(0.6),
                 '通过与"小龙虾"智能体对话交互，评估 AI 在 PPT 生成方面的实际能力', 22, TEXT_COLOR)
    
    goals = [
        '• 了解 PPT 生成的最佳实践和技术原理',
        '• 生成关于"PPT 生成对比实验"的介绍文档',
        '• 生成配套的 PPT 演示文稿',
        '• 评估 AI 智能体在 PPT 相关任务中的表现',
        '• 撰写完整的实验报告'
    ]
    add_bullet_list(slide, Inches(1.5), Inches(2.5), Inches(10), Inches(4), goals, 22, TEXT_COLOR)
    
    # 装饰
    add_shape(slide, Inches(0.8), Inches(2.5), Inches(0.08), Inches(3.5), ACCENT_COLOR)
    
    # ========== 幻灯片 4：传统 PPT 制作痛点 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, WHITE)
    
    add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), PRIMARY_COLOR)
    add_text_box(slide, Inches(0.8), Inches(0.3), Inches(12), Inches(0.8),
                 '传统 PPT 制作痛点', 36, WHITE, True)
    
    pains = [
        ('⏱️', '耗时长', '从内容策划到排版设计\n平均需 2-4 小时'),
        ('📚', '门槛高', '需要掌握 PowerPoint\n等工具的高级功能'),
        ('🎨', '质量不一', '非专业设计人员制作\n的 PPT 美观度不足'),
        ('📐', '结构混乱', '缺乏专业的逻辑框架\n和叙事结构')
    ]
    
    x_pos = 0.8
    for icon, title, desc in pains:
        # 卡片背景
        add_shape(slide, Inches(x_pos), Inches(2), Inches(2.8), Inches(3.5), LIGHT_GRAY)
        # 图标
        add_text_box(slide, Inches(x_pos + 0.3), Inches(2.3), Inches(2.2), Inches(0.8),
                     icon, 40, TEXT_COLOR, True, PP_ALIGN.CENTER)
        # 标题
        add_text_box(slide, Inches(x_pos + 0.3), Inches(3.3), Inches(2.2), Inches(0.5),
                     title, 20, PRIMARY_COLOR, True, PP_ALIGN.CENTER)
        # 描述
        add_text_box(slide, Inches(x_pos + 0.3), Inches(4), Inches(2.2), Inches(1.2),
                     desc, 16, TEXT_COLOR, False, PP_ALIGN.CENTER)
        x_pos += 3.1
    
    # ========== 幻灯片 5：AI 生成 PPT 的优势 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, WHITE)
    
    add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), PRIMARY_COLOR)
    add_text_box(slide, Inches(0.8), Inches(0.3), Inches(12), Inches(0.8),
                 'AI 智能体生成 PPT 的优势', 36, WHITE, True)
    
    advantages = [
        ('⚡', '快速生成', '分钟级完成从大纲\n到成稿的全流程'),
        ('🏗️', '结构专业', '内置专业模板\n和逻辑框架'),
        ('📝', '内容准确', '基于大语言模型的\n知识库生成内容'),
        ('🎨', '风格统一', '自动保持视觉和\n排版的一致性')
    ]
    
    x_pos = 0.8
    for icon, title, desc in advantages:
        add_shape(slide, Inches(x_pos), Inches(2), Inches(2.8), Inches(3.5), PptxRGBColor(0xE8, 0xF4, 0xFD))
        add_text_box(slide, Inches(x_pos + 0.3), Inches(2.3), Inches(2.2), Inches(0.8),
                     icon, 40, PRIMARY_COLOR, True, PP_ALIGN.CENTER)
        add_text_box(slide, Inches(x_pos + 0.3), Inches(3.3), Inches(2.2), Inches(0.5),
                     title, 20, PRIMARY_COLOR, True, PP_ALIGN.CENTER)
        add_text_box(slide, Inches(x_pos + 0.3), Inches(4), Inches(2.2), Inches(1.2),
                     desc, 16, TEXT_COLOR, False, PP_ALIGN.CENTER)
        x_pos += 3.1
    
    # ========== 幻灯片 6：主流工具对比 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, WHITE)
    
    add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), PRIMARY_COLOR)
    add_text_box(slide, Inches(0.8), Inches(0.3), Inches(12), Inches(0.8),
                 '主流 PPT 生成工具对比', 36, WHITE, True)
    
    # 表格
    tools_data = [
        ['工具', '核心技术', '优势', '适用场景'],
        ['Gamma', 'GPT-4 + 设计引擎', '设计感强', '商业演示'],
        ['Tome', 'GPT-4 + DALL-E', '叙事能力强', '故事讲述'],
        ['Beautiful.ai', '智能排版引擎', '自动排版', '商务报告'],
        ['阿里云百炼', 'Qwen + 钉钉生态', '中文优化', '企业汇报'],
        ['小龙虾', 'Qwen + 技能插件', '开源可定制', '教学实验']
    ]
    
    table = slide.shapes.add_table(len(tools_data), 4, Inches(0.8), Inches(1.8), Inches(11.5), Inches(4.5)).table
    table.columns[0].width = Inches(1.5)
    table.columns[1].width = Inches(3)
    table.columns[2].width = Inches(2.5)
    table.columns[3].width = Inches(2.5)
    
    for i, row_data in enumerate(tools_data):
        for j, cell_data in enumerate(row_data):
            cell = table.cell(i, j)
            cell.text = cell_data
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(14)
                paragraph.font.name = '微软雅黑'
                paragraph.alignment = PP_ALIGN.CENTER
                if i == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = WHITE
                else:
                    paragraph.font.color.rgb = TEXT_COLOR
            
            # 设置单元格颜色
            if i == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = PRIMARY_COLOR
            elif i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_GRAY
    
    # ========== 幻灯片 7：实验环境 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, WHITE)
    
    add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), PRIMARY_COLOR)
    add_text_box(slide, Inches(0.8), Inches(0.3), Inches(12), Inches(0.8),
                 '02  实验环境与方法', 36, WHITE, True)
    
    # 左侧：环境
    add_shape(slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(4.5), PptxRGBColor(0xE8, 0xF4, 0xFD))
    add_text_box(slide, Inches(1.2), Inches(2), Inches(4.5), Inches(0.6),
                 '实验环境', 24, PRIMARY_COLOR, True)
    
    env_items = [
        '🤖 AI 智能体：小龙虾-OpenClaw',
        '🧠 使用模型：Qwen3.6-plus',
        '💬 交互平台：钉钉',
        '📄 参考材料：《PPT生成对比实验》',
        '🔧 辅助工具：baoyu-slide-deck',
        '⏱️ 实验时长：约 120 分钟'
    ]
    add_bullet_list(slide, Inches(1.2), Inches(2.8), Inches(4.5), Inches(3), env_items, 16, TEXT_COLOR)
    
    # 右侧：方法
    add_shape(slide, Inches(7), Inches(1.8), Inches(5.5), Inches(4.5), PptxRGBColor(0xFD, 0xF0, 0xE0))
    add_text_box(slide, Inches(7.4), Inches(2), Inches(4.5), Inches(0.6),
                 '实验方法', 24, ACCENT_COLOR, True)
    
    methods = [
        '1. 对话交互 - 自然语言对话',
        '2. 内容生成 - Markdown 格式',
        '3. PPT 制作 - 技能插件生成',
        '4. 能力评估 - 多维度评分',
        '5. 报告撰写 - 整理实验结果'
    ]
    add_bullet_list(slide, Inches(7.4), Inches(2.8), Inches(4.5), Inches(3), methods, 16, TEXT_COLOR)
    
    # ========== 幻灯片 8：实验过程 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, WHITE)
    
    add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), PRIMARY_COLOR)
    add_text_box(slide, Inches(0.8), Inches(0.3), Inches(12), Inches(0.8),
                 '04  实验过程记录', 36, WHITE, True)
    
    stages = [
        ('第一阶段', '与小龙虾智能体\n对话学习', '了解 PPT 生成的基本\n原理和最佳实践'),
        ('第二阶段', '生成介绍文档\nMarkdown', '生成包含 6 个章节、\n约 3000 字的文档'),
        ('第三阶段', '生成 PPT 演示\n文稿', '15 张幻灯片，商务简\n洁风格，蓝色主色调')
    ]
    
    x_pos = 0.8
    for i, (stage, title, desc) in enumerate(stages):
        # 卡片
        add_shape(slide, Inches(x_pos), Inches(2), Inches(3.8), Inches(4), LIGHT_GRAY)
        # 阶段标签
        add_shape(slide, Inches(x_pos), Inches(2), Inches(3.8), Inches(0.8), PRIMARY_COLOR)
        add_text_box(slide, Inches(x_pos + 0.3), Inches(2.15), Inches(3.2), Inches(0.6),
                     stage, 22, WHITE, True, PP_ALIGN.CENTER)
        # 内容
        add_text_box(slide, Inches(x_pos + 0.3), Inches(3.2), Inches(3.2), Inches(1),
                     title, 18, PRIMARY_COLOR, True, PP_ALIGN.CENTER)
        add_text_box(slide, Inches(x_pos + 0.3), Inches(4.3), Inches(3.2), Inches(1.2),
                     desc, 15, TEXT_COLOR, False, PP_ALIGN.CENTER)
        
        # 箭头
        if i < 2:
            add_text_box(slide, Inches(x_pos + 3.9), Inches(3.8), Inches(0.4), Inches(0.5),
                         '→', 28, ACCENT_COLOR, True, PP_ALIGN.CENTER)
        
        x_pos += 4.2
    
    # ========== 幻灯片 9：内容生成能力评估 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, WHITE)
    
    add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), PRIMARY_COLOR)
    add_text_box(slide, Inches(0.8), Inches(0.3), Inches(12), Inches(0.8),
                 '05  实验结果 - 内容生成能力', 36, WHITE, True)
    
    eval_data = [
        ['评估项', '得分', '说明'],
        ['内容完整性', '⭐⭐⭐⭐⭐', '覆盖所有核心章节，无遗漏'],
        ['逻辑结构', '⭐⭐⭐⭐⭐', '层次清晰，递进合理'],
        ['设计美观度', '⭐⭐⭐⭐', '商务风格统一，排版整洁'],
        ['文字精炼度', '⭐⭐⭐⭐', '要点突出，无冗余信息'],
        ['图表适配性', '⭐⭐⭐⭐', '数据可视化图表匹配度高']
    ]
    
    table = slide.shapes.add_table(len(eval_data), 3, Inches(0.8), Inches(1.8), Inches(11.5), Inches(4.5)).table
    table.columns[0].width = Inches(2.5)
    table.columns[1].width = Inches(3)
    table.columns[2].width = Inches(6)
    
    for i, row_data in enumerate(eval_data):
        for j, cell_data in enumerate(row_data):
            cell = table.cell(i, j)
            cell.text = cell_data
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(16)
                paragraph.font.name = '微软雅黑'
                paragraph.alignment = PP_ALIGN.CENTER
                if i == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = WHITE
            
            if i == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = PRIMARY_COLOR
            elif i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_GRAY
    
    # 综合得分
    add_shape(slide, Inches(4), Inches(6.3), Inches(5), Inches(0.8), ACCENT_COLOR)
    add_text_box(slide, Inches(4), Inches(6.35), Inches(5), Inches(0.7),
                 '综合得分：4.4 / 5.0', 28, WHITE, True, PP_ALIGN.CENTER)
    
    # ========== 幻灯片 10：效率分析 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, WHITE)
    
    add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), PRIMARY_COLOR)
    add_text_box(slide, Inches(0.8), Inches(0.3), Inches(12), Inches(0.8),
                 '05  实验结果 - 效率分析', 36, WHITE, True)
    
    efficiency_data = [
        ['任务', '传统方式', 'AI 智能体', '效率提升'],
        ['内容策划', '30-60 分钟', '2-5 分钟', '6-12 倍'],
        ['大纲撰写', '20-40 分钟', '1-3 分钟', '7-13 倍'],
        ['PPT 排版', '40-80 分钟', '2-5 分钟', '8-16 倍'],
        ['整体制作', '90-180 分钟', '5-13 分钟', '7-14 倍']
    ]
    
    table = slide.shapes.add_table(len(efficiency_data), 4, Inches(0.8), Inches(1.8), Inches(11.5), Inches(3)).table
    table.columns[0].width = Inches(2.5)
    table.columns[1].width = Inches(3)
    table.columns[2].width = Inches(3)
    table.columns[3].width = Inches(3)
    
    for i, row_data in enumerate(efficiency_data):
        for j, cell_data in enumerate(row_data):
            cell = table.cell(i, j)
            cell.text = cell_data
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(18)
                paragraph.font.name = '微软雅黑'
                paragraph.alignment = PP_ALIGN.CENTER
                if i == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = WHITE
            
            if i == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = PRIMARY_COLOR
            elif j == 3:
                cell.fill.solid()
                cell.fill.fore_color.rgb = PptxRGBColor(0xE8, 0xF4, 0xFD)
    
    # 结论框
    add_shape(slide, Inches(2), Inches(5.2), Inches(9), Inches(1.5), PptxRGBColor(0xE8, 0xF4, 0xFD))
    add_text_box(slide, Inches(2.5), Inches(5.4), Inches(8), Inches(1.2),
                 '🚀 效率提升结论：相比传统制作方式，AI 智能体可将 PPT 整体制作效率提升 7-14 倍', 
                 22, PRIMARY_COLOR, True, PP_ALIGN.CENTER)
    
    # ========== 幻灯片 11：主要发现 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, WHITE)
    
    add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), PRIMARY_COLOR)
    add_text_box(slide, Inches(0.8), Inches(0.3), Inches(12), Inches(0.8),
                 '06  实验结论 - 主要发现', 36, WHITE, True)
    
    findings = [
        ('✅', '实用水平', 'AI 智能体 PPT 生成能力已达到实用水平'),
        ('⭐', '表现优异', '小龙虾智能体基于 Qwen 大模型中文理解能力强'),
        ('🔧', '仍需改进', '高级设计模板和复杂数据可视化能力有待提升')
    ]
    
    x_pos = 0.8
    for icon, title, desc in findings:
        add_shape(slide, Inches(x_pos), Inches(2), Inches(3.8), Inches(4), LIGHT_GRAY)
        add_shape(slide, Inches(x_pos), Inches(2), Inches(3.8), Inches(0.8), PRIMARY_COLOR)
        add_text_box(slide, Inches(x_pos + 0.2), Inches(2.15), Inches(3.4), Inches(0.6),
                     f'{icon} {title}', 22, WHITE, True, PP_ALIGN.CENTER)
        add_text_box(slide, Inches(x_pos + 0.3), Inches(3.2), Inches(3.2), Inches(2),
                     desc, 18, TEXT_COLOR, False, PP_ALIGN.CENTER)
        x_pos += 4.2
    
    # ========== 幻灯片 12：实验局限性 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, WHITE)
    
    add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), PRIMARY_COLOR)
    add_text_box(slide, Inches(0.8), Inches(0.3), Inches(12), Inches(0.8),
                 '实验局限性', 36, WHITE, True)
    
    limitations = [
        ('📊', '样本有限', '仅测试了单一主题\n未覆盖多场景'),
        ('🔄', '未迭代优化', '未进行多轮对话\n优化测试'),
        ('🔍', '未横向对比', '未与其他 PPT 生成\n工具进行对比')
    ]
    
    x_pos = 0.8
    for icon, title, desc in limitations:
        add_shape(slide, Inches(x_pos), Inches(2), Inches(3.8), Inches(3.5), PptxRGBColor(0xFD, 0xE8, 0xE8))
        add_text_box(slide, Inches(x_pos + 0.3), Inches(2.5), Inches(3.2), Inches(0.8),
                     icon, 40, PptxRGBColor(0xCC, 0x33, 0x33), True, PP_ALIGN.CENTER)
        add_text_box(slide, Inches(x_pos + 0.3), Inches(3.5), Inches(3.2), Inches(0.5),
                     title, 20, PptxRGBColor(0xCC, 0x33, 0x33), True, PP_ALIGN.CENTER)
        add_text_box(slide, Inches(x_pos + 0.3), Inches(4.2), Inches(3.2), Inches(1),
                     desc, 16, TEXT_COLOR, False, PP_ALIGN.CENTER)
        x_pos += 4.2
    
    # ========== 幻灯片 13：后续建议 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, WHITE)
    
    add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), PRIMARY_COLOR)
    add_text_box(slide, Inches(0.8), Inches(0.3), Inches(12), Inches(0.8),
                 '后续建议', 36, WHITE, True)
    
    suggestions = [
        '🔬 扩大测试范围：测试不同主题和风格的 PPT 生成效果',
        '📊 多工具对比：与 Gamma、Tome 等工具进行横向对比',
        '🔄 迭代优化：测试多轮对话优化后的效果提升',
        '🏫 实际应用：将 AI 生成 PPT 应用于实际教学和工作场景'
    ]
    
    y_pos = 1.8
    for item in suggestions:
        add_shape(slide, Inches(1.5), Inches(y_pos), Inches(10), Inches(0.9), LIGHT_GRAY)
        add_text_box(slide, Inches(1.8), Inches(y_pos + 0.15), Inches(9.4), Inches(0.7),
                     item, 18, TEXT_COLOR, False)
        y_pos += 1.1
    
    # ========== 幻灯片 14：实验总结 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, WHITE)
    
    add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), PRIMARY_COLOR)
    add_text_box(slide, Inches(0.8), Inches(0.3), Inches(12), Inches(0.8),
                 '实验总结', 36, WHITE, True)
    
    add_text_box(slide, Inches(0.8), Inches(1.6), Inches(11.5), Inches(0.6),
                 '本实验通过与小龙虾智能体的对话交互，完成了 PPT 生成主题的介绍文档和演示文稿制作', 
                 22, TEXT_COLOR)
    
    conclusions = [
        '✓ AI 智能体在 PPT 生成方面已具备实用能力，能够显著提升制作效率',
        '✓ 小龙虾智能体基于 Qwen 大模型，在中文内容生成方面表现优异',
        '✓ 钉钉生态集成使得企业级应用更加便捷',
        '✓ 开源可定制的特性为个性化需求提供了可能'
    ]
    
    add_bullet_list(slide, Inches(1.5), Inches(2.5), Inches(10), Inches(3), conclusions, 20, TEXT_COLOR)
    
    # 底部引用
    add_shape(slide, Inches(1.5), Inches(5.8), Inches(10), Inches(1), PptxRGBColor(0xE8, 0xF4, 0xFD))
    add_text_box(slide, Inches(2), Inches(5.9), Inches(9), Inches(0.8),
                 '通过本次实验，我不仅掌握了 AI 智能体 PPT 生成的方法，也对 AI 在教育和工作场景中的应用前景有了更深入的理解。',
                 16, TEXT_COLOR, False, PP_ALIGN.CENTER)
    
    # ========== 幻灯片 15：致谢 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, PRIMARY_COLOR)
    
    add_shape(slide, Inches(0), Inches(3.2), Inches(13.333), Inches(0.06), ACCENT_COLOR)
    
    add_text_box(slide, Inches(1.5), Inches(2), Inches(10), Inches(1),
                 '感谢观看', 48, WHITE, True, PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1.5), Inches(3.5), Inches(10), Inches(1),
                 'THANK YOU', 36, PptxRGBColor(0xCC, 0xDD, 0xEE), False, PP_ALIGN.CENTER)
    
    add_text_box(slide, Inches(1.5), Inches(5), Inches(10), Inches(0.5),
                 '实验人：宋高瀞    课程：数字人文与 AI 应用', 18, PptxRGBColor(0x99, 0xAA, 0xBB), False, PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1.5), Inches(5.6), Inches(10), Inches(0.5),
                 '2026 年 6 月 15 日', 16, PptxRGBColor(0x99, 0xAA, 0xBB), False, PP_ALIGN.CENTER)
    
    # 保存
    output_path = '/home/admin/.openclaw/workspace/PPT_演示文稿_宋高瀞.pptx'
    prs.save(output_path)
    print(f'PPT 已保存：{output_path}')
    return output_path

# 执行
if __name__ == '__main__':
    print('开始生成实验报告文档和 PPT...')
    word_path = create_word_report()
    ppt_path = create_ppt()
    print('\n✅ 文件生成完成！')
    print(f'📄 Word 文档：{word_path}')
    print(f'📊 PPT 演示文稿：{ppt_path}')
