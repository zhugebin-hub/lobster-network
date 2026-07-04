#!/usr/bin/env python3
"""
"101计划"首批核心课程培育推进会汇报PPT
基于《基于AI原生教学范式的计算机网络课程改革实践》论文
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

C = {
    'primary': RGBColor(0x1A, 0x36, 0x5D),
    'secondary': RGBColor(0x2B, 0x6C, 0xB0),
    'accent': RGBColor(0xE8, 0x83, 0x2A),
    'light': RGBColor(0xF7, 0xF9, 0xFC),
    'white': RGBColor(0xFF, 0xFF, 0xFF),
    'text': RGBColor(0x2D, 0x37, 0x48),
    'muted': RGBColor(0x71, 0x80, 0x96),
    'success': RGBColor(0x38, 0xA1, 0x69),
    'teal': RGBColor(0x31, 0x97, 0x95),
    'purple': RGBColor(0x9F, 0x7A, 0xEA),
}
FN = 'Microsoft YaHei'

def bg(s, c):
    f = s.background.fill; f.solid(); f.fore_color.rgb = c

def rect(s, x, y, w, h, fc, lc=None, lw=0, r=False):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if r else MSO_SHAPE.RECTANGLE,
                             Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fc
    if lc: sh.line.color.rgb = lc; sh.line.width = Pt(lw)
    else: sh.line.fill.background()
    return sh

def tb(s, x, y, w, h, t, fs=14, c=C['text'], b=False, a=PP_ALIGN.LEFT, italic=False):
    bx = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = bx.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = t; p.font.size = Pt(fs); p.font.color.rgb = c
    p.font.bold = b; p.font.italic = italic; p.font.name = FN; p.alignment = a
    return bx

def mt(s, x, y, w, h, lines, fs=14, c=C['text'], b=False, a=PP_ALIGN.LEFT, italic=False, sp=2):
    bx = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = bx.text_frame; tf.word_wrap = True
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = line; p.font.size = Pt(fs); p.font.color.rgb = c
        p.font.bold = b; p.font.italic = italic; p.font.name = FN; p.alignment = a
        p.space_after = Pt(sp)
    return bx

def hdr(s, title):
    rect(s, 0, 0, 13.33, 1.1, C['primary'])
    tb(s, 0.5, 0.15, 12, 0.9, title, fs=26, c=C['white'], b=True)

# ============================================================
# SLIDE 1: Cover
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, C['primary'])
rect(s, 0, 0, 13.33, 0.12, C['accent'])
tb(s, 1, 0.8, 11.33, 1.0,
   '\u6559\u80b2\u90e8\u201c101\u8ba1\u5212\u201d\u9996\u6279\u6838\u5fc3\u8bfe\u7a0b\u57f9\u80b2\u63a8\u8fdb\u4f1a',
   fs=34, c=C['white'], b=True, a=PP_ALIGN.CENTER)
tb(s, 1, 2.1, 11.33, 0.8,
   '\u57fa\u4e8eAI\u539f\u751f\u6559\u5b66\u8303\u5f0f\u7684\u8ba1\u7b97\u673a\u7f51\u7edc\u8bfe\u7a0b\u6539\u9769\u5b9e\u8df5',
   fs=26, c=C['accent'], b=True, a=PP_ALIGN.CENTER)
tb(s, 1, 3.3, 11.33, 0.5,
   '\u2014\u2014 \u5f15\u81ea\u300a\u9ad8\u7b49\u5de5\u7a0b\u6559\u80b2\u7814\u7a76\u300b\u6295\u7a3f\u8bba\u6587',
   fs=16, c=C['muted'], italic=True, a=PP_ALIGN.CENTER)
mt(s, 1, 5.5, 11.33, 0.8,
   ['\u6d59\u6c5f\u5de5\u5546\u5927\u5b66 \u00b7 \u8bf8\u845b\u658c\u56e2\u961f',
    '\u8bba\u6587\u4f5c\u8005\uff1a\u8bf8\u845b\u658c\u3001\u91d1\u84c9\u3001\u9ad8\u660e\u3001\u674e\u4f20\u714c\u3001\u848b\u732e',
    '2026\u5e746\u670825\u65e5 | \u79d1\u521b\u5927\u697c206'],
   fs=14, c=C['muted'], a=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 2: Abstract
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, C['white'])
hdr(s, '\u7814\u7a76\u6458\u8981')

rect(s, 0.8, 1.3, 11.5, 2.2, C['light'], C['accent'], 2, True)
tb(s, 1, 1.4, 2, 0.5, '\u6458\u8981', fs=18, c=C['primary'], b=True)
mt(s, 1, 1.9, 11, 1.5,
   ['\u672c\u6587\u4ee5\u6d59\u6c5f\u7701\u201cCS&AI 101\u8ba1\u5212\u201d\u8ba1\u7b97\u673a\u7f51\u7edc\u8bfe\u7a0b\u5efa\u8bbe\u4e3a\u80cc\u666f\uff0c\u63a2\u8ba8\u5982\u4f55\u5c06AI\u5de5\u5177\u878d\u5165\u6559\u5b66\u5b9e\u8df5\uff0c\u5b9e\u73b0\u4ece\u201c\u77e5\u8bc6\u704c\u8f93\u201d\u5411\u201cAI\u5e94\u7528\u5b9e\u8df5\u201d\u7684\u8303\u5f0f\u8f6c\u578b\u3002\u8bfe\u7a0b\u56e2\u961f\u63d0\u51fa\u4e86\u201c\u8bfe\u7a0b\u4f5c\u4e3aAI\u5e94\u7528\u80cc\u666f\u573a\u666f\u201d\u7684\u6838\u5fc3\u7406\u5ff5\uff0c\u6784\u5efa\u4e86\u201c\u5b66\u751f\u4f5c\u4e1a\u5373\u8bfe\u7a0b\u8d44\u6e90\u201d\u7684\u521b\u65b0\u6a21\u5f0f\uff0c\u5e76\u901a\u8fc7SDP\u6848\u4f8b\u5b9e\u8df5\u9a8c\u8bc1\u4e86AI\u8fed\u4ee3\u601d\u7ef4\u5728\u6559\u5b66\u4e2d\u7684\u53ef\u884c\u6027\u3002'],
   fs=14, c=C['text'])

tb(s, 0.8, 3.8, 5, 0.5, '\u5173\u952e\u8bcd', fs=18, c=C['primary'], b=True)
mt(s, 0.8, 4.3, 11, 1.0,
   ['AI\u539f\u751f\u6559\u5b66 | \u8ba1\u7b97\u673a\u7f51\u7edc | \u8bfe\u7a0b\u5efa\u8bbe | 101\u8ba1\u5212 | \u6559\u80b2\u8303\u5f0f\u91cd\u6784 | \u667a\u80fd\u4f53\u8f85\u52a9\u6559\u5b66'],
   fs=15, c=C['accent'])

tb(s, 0.8, 5.3, 5, 0.5, '\u73b0\u6709\u57fa\u7840\u6570\u636e', fs=18, c=C['primary'], b=True)
mt(s, 0.8, 5.8, 11, 1.2,
   ['MOOC\u5728\u7ebf\u8bfe\u7a0b\uff1a228\u4eba\u5b66\u4e60 | \u7701\u7ea7\u7ebf\u4e0a\u4e00\u6d41\u5b9e\u9a8c\u8bfe\uff1a7\u671f/623\u4eba\u6b21/22\u6821/59.4\u4e07\u8bbf\u95ee',
    '\u5df2\u51fa\u7248\u6559\u6750\u4e24\u90e8\uff08\u6e05\u534e\u5927\u5b66\u51fa\u7248\u793e\uff09'],
   fs=14, c=C['text'])

# ============================================================
# SLIDE 3: Research Background
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, C['white'])
hdr(s, '\u4e00\u3001\u7814\u7a76\u80cc\u666f\u4e0e\u6311\u6218')

tb(s, 0.8, 1.4, 5, 0.5, '101\u8ba1\u5212\u80cc\u666f', fs=18, c=C['primary'], b=True)
mt(s, 0.8, 1.9, 5.5, 2.5,
   ['\u2022 2026\u5e74\u6559\u80b2\u90e8\u201c101\u8ba1\u5212\u201d\u5168\u9762\u542f\u52a8',
    '\u2022 \u901a\u8fc7\u6838\u5fc3\u8bfe\u7a0b\u3001\u6838\u5fc3\u6559\u6750\u3001\u6838\u5fc3\u5b9e\u8df5\u9879\u76ee\u548c\u6838\u5fc3\u5e08\u8d44\u56e2\u961f\u5efa\u8bbe\u63a8\u52a8\u9ad8\u7b49\u6559\u80b2\u8d28\u91cf\u63d0\u5347',
    '\u2022 \u6d59\u6c5f\u7701\u201cCS&AI 101\u8ba1\u5212\u201d\u805a\u7126\u8ba1\u7b97\u673a\u79d1\u5b66\u4e0e\u4eba\u5de5\u667a\u80fd\u6df1\u5ea6\u878d\u5408',
    '\u2022 \u8ba1\u7b97\u673a\u7f51\u7edc\u4f5c\u4e3a\u8ba1\u7b97\u673a\u7c7b\u4e13\u4e1a\u6838\u5fc3\u57fa\u7840\u8bfe\u7a0b'],
   fs=14, c=C['text'])

tb(s, 7, 1.4, 5, 0.5, '\u4f20\u7edf\u6559\u5b66\u6a21\u5f0f\u9762\u4e34\u7684\u6311\u6218', fs=18, c=C['accent'], b=True)
mt(s, 7, 1.9, 5.5, 2.5,
   ['1. \u77e5\u8bc6\u66f4\u65b0\u6ede\u540e\uff1a\u7f51\u7edc\u6280\u672f\u53d1\u5c55\u8fc5\u901f\uff0c\u6559\u6750\u5185\u5bb9\u96be\u4ee5\u8ddf\u4e0a\u884c\u4e1a\u53d8\u5316',
    '2. \u5b9e\u8df5\u73af\u8282\u8584\u5f31\uff1a\u5b9e\u9a8c\u8bbe\u5907\u6210\u672c\u9ad8\uff0c\u5b66\u751f\u52a8\u624b\u673a\u4f1a\u6709\u9650',
    '3. \u6559\u5b66\u6548\u679c\u8bc4\u4f30\u5355\u4e00\uff1a\u4f20\u7edf\u4f5c\u4e1a\u5f62\u5f0f\u96be\u4ee5\u53cd\u6620\u5b66\u751f\u771f\u5b9e\u80fd\u529b',
    '4. \u6559\u5b66\u8d44\u6e90\u5206\u6563\uff1a\u4f18\u8d28\u8d44\u6e90\u96be\u4ee5\u7cfb\u7edf\u5316\u79ef\u7d2f\u548c\u5171\u4eab'],
   fs=14, c=C['text'])

rect(s, 0.8, 4.6, 11.5, 2.2, C['primary'], r=True)
tb(s, 1, 4.7, 3, 0.5, '\u7814\u7a76\u610f\u4e49', fs=18, c=C['accent'], b=True)
mt(s, 1, 5.2, 11, 1.5,
   ['\u2022 \u7406\u8bba\u5c42\u9762\uff1a\u63d0\u51fa\u201cAI\u539f\u751f\u6559\u5b66\u8303\u5f0f\u201d\u6982\u5ff5\uff0c\u91cd\u65b0\u5b9a\u4e49\u8bfe\u7a0b\u3001\u6559\u5e08\u3001\u5b66\u751f\u5728AI\u65f6\u4ee3\u7684\u89d2\u8272\u5b9a\u4f4d',
    '\u2022 \u5b9e\u8df5\u5c42\u9762\uff1a\u6784\u5efa\u201c\u5b66\u751f\u4f5c\u4e1a\u5373\u8bfe\u7a0b\u8d44\u6e90\u201d\u7684\u521b\u65b0\u6a21\u5f0f\uff0c\u5b9e\u73b0\u6559\u5b66\u8d44\u6e90\u7684AI\u539f\u751f\u91cd\u6784',
    '\u2022 \u63a8\u5e7f\u5c42\u9762\uff1a\u5f62\u6210\u53ef\u590d\u5236\u3001\u53ef\u63a8\u5e7f\u7684\u201cAI+\u6559\u80b2\u201d\u5178\u578b\u6848\u4f8b\uff0c\u4e3a\u5176\u4ed6\u8bfe\u7a0b\u5efa\u8bbe\u63d0\u4f9b\u53c2\u8003'],
   fs=14, c=C['white'])

# ============================================================
# SLIDE 4: Literature Review
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, C['white'])
hdr(s, '\u4e8c\u3001AI\u6559\u80b2\u7814\u7a76\u73b0\u72b6\u4e0e\u7406\u8bba\u6846\u67b6')

tb(s, 0.8, 1.4, 5, 0.5, 'AI\u5728\u6559\u80b2\u9886\u57df\u7684\u5e94\u7528\u7814\u7a76', fs=18, c=C['primary'], b=True)
mt(s, 0.8, 1.9, 5.5, 2.5,
   ['\u2022 AI\u8f85\u52a9\u6559\u5b66\uff1aAI\u8f85\u52a9\u6559\u5e08\u8fdb\u884c\u6559\u5b66\u8bbe\u8ba1\u3001\u5185\u5bb9\u751f\u6210\u3001\u4f5c\u4e1a\u6279\u6539\uff08Smith et al., 2023\uff09',
    '\u2022 \u4e2a\u6027\u5316\u5b66\u4e60\uff1a\u57fa\u4e8eAI\u7684\u5b66\u4e60\u5206\u6790\u6280\u672f\uff0c\u63d0\u4f9b\u4e2a\u6027\u5316\u5b66\u4e60\u8def\u5f84\u548c\u8d44\u6e90\u63a8\u8350\uff08Chen & Wang, 2024\uff09',
    '\u2022 \u667a\u80fd\u8bc4\u4f30\uff1a\u5229\u7528AI\u6280\u672f\u8fdb\u884c\u5b66\u4e60\u6548\u679c\u8bc4\u4f30\uff0c\u5b9e\u73b0\u81ea\u52a8\u5316\u8bc4\u5206\u548c\u53cd\u9988\uff08Li et al., 2025\uff09'],
   fs=13, c=C['text'])

tb(s, 7, 1.4, 5, 0.5, '\u5de5\u7a0b\u6559\u80b2\u8303\u5f0f\u8f6c\u578b', fs=18, c=C['primary'], b=True)
mt(s, 7, 1.9, 5.5, 2.5,
   ['\u2022 CDIO\u6a21\u5f0f\uff1a\u5f3a\u8c03\u201c\u6784\u601d-\u8bbe\u8ba1-\u5b9e\u73b0-\u8fd0\u4f5c\u201d\u7684\u5de5\u7a0b\u6559\u80b2\u6a21\u5f0f',
    '\u2022 \u65b0\u5de5\u79d1\u5efa\u8bbe\uff1a\u63a2\u7d22AI\u65f6\u4ee3\u5de5\u7a0b\u6559\u80b2\u7684\u65b0\u6a21\u5f0f\u3001\u65b0\u8def\u5f84\uff08\u674e\u57f9\u6839\u7b49\uff0c2023\uff09',
    '\u2022 \u4ea7\u6559\u878d\u5408\uff1a\u63a8\u52a8\u9ad8\u6821\u4e0e\u4f01\u4e1a\u6df1\u5ea6\u5408\u4f5c\uff08\u738b\u5efa\u534e\uff0c2024\uff09'],
   fs=13, c=C['text'])

rect(s, 0.8, 4.6, 11.5, 2.4, C['light'], C['secondary'], 2, True)
tb(s, 1, 4.7, 11, 0.5, 'AI\u539f\u751f\u6559\u5b66\u8303\u5f0f\u7406\u8bba\u6846\u67b6', fs=18, c=C['primary'], b=True, a=PP_ALIGN.CENTER)

rect(s, 1.2, 5.2, 3.3, 1.6, C['secondary'], r=True)
tb(s, 1.3, 5.25, 3, 0.3, '\u8bfe\u7a0b\u4f5c\u4e3a\u573a\u666f', fs=14, c=C['white'], b=True, a=PP_ALIGN.CENTER)
tb(s, 1.3, 5.6, 3, 1.0, '\u8bfe\u7a0b\u5185\u5bb9\u4e0d\u518d\u662f\u76ee\u7684\uff0c\u800c\u662f\u8bad\u7ec3AI\u5de5\u5177\u5e94\u7528\u80fd\u529b\u7684\u8f7d\u4f53', fs=11, c=C['white'], a=PP_ALIGN.CENTER)

rect(s, 4.9, 5.2, 3.3, 1.6, C['teal'], r=True)
tb(s, 5, 5.25, 3, 0.3, 'AI\u4f5c\u4e3a\u5de5\u5177', fs=14, c=C['white'], b=True, a=PP_ALIGN.CENTER)
tb(s, 5, 5.6, 3, 1.0, '\u5b66\u751f\u4f7f\u7528AI\u5de5\u5177\u5b8c\u6210\u77e5\u8bc6\u5b66\u4e60\u3001\u5185\u5bb9\u751f\u6210\u3001\u95ee\u9898\u89e3\u51b3', fs=11, c=C['white'], a=PP_ALIGN.CENTER)

rect(s, 8.6, 5.2, 3.3, 1.6, C['accent'], r=True)
tb(s, 8.7, 5.25, 3, 0.3, '\u6559\u5e08\u4f5c\u4e3a\u5f15\u5bfc\u8005', fs=14, c=C['white'], b=True, a=PP_ALIGN.CENTER)
tb(s, 8.7, 5.6, 3, 1.0, '\u4ece\u5185\u5bb9\u4f20\u6388\u8005\u8f6c\u53d8\u4e3aAI\u5e94\u7528\u7684\u5f15\u5bfc\u8005\u4e0e\u8d28\u91cf\u628a\u5173\u8005', fs=11, c=C['white'], a=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 5: Current Course Foundation
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, C['white'])
hdr(s, '\u4e09\u3001\u73b0\u6709\u8bfe\u7a0b\u57fa\u7840')

ach = [
    { 't': '\u7406\u8bba\u8bfe\u7a0b\uff1aMOOC',
      'd': '\u300a\u9ad8\u7ea7\u7f51\u7edc\u901a\u4fe1\u539f\u7406\u5b9e\u8df5\u300b\n\u4e2d\u56fd\u5927\u5b66MOOC\u5e73\u53f0 | 2026\u6625\u5b63\n18\u5468/13\u7ae0\uff08\u5df2\u5b8c\u621014\u5468\uff09\n\u5b66\u4e60\u4eba\u6570\uff1a228\u4eba | \u4e92\u52a8\uff1a2950\u6b21\n\u5408\u4f5c\u4f01\u4e1a\uff1a\u963f\u91cc\u4e91', 'c': C['secondary'] },
    { 't': '\u5b9e\u9a8c\u8bfe\u7a0b\uff1a\u7701\u7ea7\u4e00\u6d41',
      'd': '\u300a\u8ba1\u7b97\u673a\u7f51\u7edc\u5b9e\u9a8c\u300b\n\u6d59\u6c5f\u7701\u9ad8\u7b49\u5b66\u6821\u5728\u7ebf\u5f00\u653e\u8bfe\u7a0b\u5e73\u53f0\n7\u671f\u8fd0\u884c | 623\u4eba\u6b21\u9009\u8bfe\n22\u6240\u9ad8\u6821\u8986\u76d6\n\u7d2f\u8ba1\u8bbf\u95ee593,924\u6b21', 'c': C['teal'] },
    { 't': '\u6559\u6750\u5efa\u8bbe\uff1a\u6e05\u534e\u51fa\u7248\u793e',
      'd': '\u300a\u7f51\u7edc\u901a\u4fe1\u539f\u7406\u5b9e\u8df5\u300b\uff082024.9\uff09\n\u5fae\u8bfe\u89c6\u9891\u7248\uff0c615\u5206\u949f\u89c6\u9891\u8bb2\u89e3\n\u300a\u7cfb\u7edf\u7ea7\u7f16\u7a0b\u53ca\u5206\u5e03\u5f0f\u5e94\u7528\u5b9e\u73b0\u6280\u672f\u300b\n\u5df2\u5b9a\u7a3f\uff0c\u9884\u8ba12026\u5e74\u521d\u51fa\u7248', 'c': C['success'] },
]

ax = 0.5
for a in ach:
    rect(s, ax, 1.4, 3.95, 4.0, C['light'], a['c'], 2, True)
    tb(s, ax+0.15, 1.5, 3.65, 0.5, a['t'], fs=17, c=a['c'], b=True, a=PP_ALIGN.CENTER)
    mt(s, ax+0.15, 2.1, 3.65, 3.0, a['d'].split('\n'), fs=12, c=C['text'])
    ax += 4.15

# Team
rect(s, 0.8, 5.6, 11.5, 1.3, C['primary'], r=True)
tb(s, 1, 5.65, 11, 0.4, '\u8bfe\u7a0b\u56e2\u961f\u4ecb\u7ecd', fs=16, c=C['accent'], b=True, a=PP_ALIGN.CENTER)
mt(s, 1, 6.05, 11, 0.8,
   ['\u8bf8\u845b\u658c\uff08\u6559\u6388\uff09\uff1a\u8bfe\u7a0b\u8d1f\u8d23\u4eba | \u91d1\u84c9\uff08\u526f\u6559\u6388\uff09\uff1a\u5b9e\u9a8c\u8bfe\u8d1f\u8d23\u4eba | \u9ad8\u660e\uff08\u526f\u6559\u6388/\u7cfb\u4e3b\u4efb\uff09\uff1a\u6559\u6750\u5efa\u8bbe',
    '\u674e\u4f20\u714c\uff08\u6559\u6388/\u90e8\u957f\uff09\uff1a\u6821\u4f01\u5408\u4f5c | \u848b\u732e\uff08\u5b9e\u9a8c\u5e08\uff09\uff1a\u6559\u5b66\u8d44\u6e90\u5f00\u53d1'],
   fs=13, c=C['white'])

# ============================================================
# SLIDE 6: Core Philosophy
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, C['white'])
hdr(s, '\u56db\u3001AI\u539f\u751f\u6559\u5b66\u8303\u5f0f\u6838\u5fc3\u7406\u5ff5')

rect(s, 0.8, 1.3, 3.5, 4.5, C['primary'], r=True)
tb(s, 1, 1.4, 3, 0.5, '\u8bfe\u7a0b\u5b9a\u4f4d\u8f6c\u578b', fs=18, c=C['accent'], b=True, a=PP_ALIGN.CENTER)
mt(s, 1, 1.9, 3, 3.5,
   ['\u4ece\u201c\u77e5\u8bc6\u704c\u8f93\u201d',
    '\u2193',
    '\u5230\u201cAI\u5e94\u7528\u5b9e\u8df5\u201d',
    '',
    '\u201c\u8bfe\u7a0b\u4e0d\u518d\u662f\u6559\u5b66\u6838\u5fc3\uff0c',
    '  \u800c\u662fAI\u5de5\u5177\u5e94\u7528\u7684\u80cc\u666f\u573a\u666f\u3002\u201d',
    '',
    '\u4e09\u5c42\u542b\u4e49\uff1a',
    '1. \u8bfe\u7a0b\u4f5c\u4e3a\u573a\u666f',
    '2. AI\u4f5c\u4e3a\u5de5\u5177',
    '3. \u6559\u5e08\u4f5c\u4e3a\u5f15\u5bfc\u8005'],
   fs=13, c=C['white'])

rect(s, 4.7, 1.3, 3.5, 4.5, C['teal'], r=True)
tb(s, 4.9, 1.4, 3, 0.5, '\u6559\u80b2\u76ee\u6807\u91cd\u6784', fs=18, c=C['white'], b=True, a=PP_ALIGN.CENTER)
mt(s, 4.9, 1.9, 3, 3.5,
   ['\u4ece\u201c\u77e5\u8bc6\u8bb0\u5fc6\u201d',
    '\u2193',
    '\u5230\u201c\u5224\u65ad\u529b\u57f9\u517b\u201d',
    '',
    '\u201c\u672a\u6765\u8bfe\u7a0b\u5efa\u8bbe\u7684\u5173\u952e\u4e0d\u518d\u662f\u77e5\u8bc6\u4f20\u6388\uff0c\u800c\u662f\u57f9\u517b\u5b66\u751f\u5bf9AI\u751f\u6210\u5185\u5bb9\u7684\u5224\u65ad\u529b\u3002\u201d',
    '',
    '\u6838\u5fc3\u7d20\u517b\uff1a',
    '1. \u5224\u65ad\u529b',
    '2. \u6279\u5224\u6027\u601d\u7ef4',
    '3. \u4eba\u673a\u534f\u540c\u51b3\u7b56'],
   fs=13, c=C['white'])

rect(s, 8.6, 1.3, 3.9, 4.5, C['accent'], r=True)
tb(s, 8.8, 1.4, 3.4, 0.5, '\u672a\u6765\u5c55\u671b', fs=18, c=C['white'], b=True, a=PP_ALIGN.CENTER)
mt(s, 8.8, 1.9, 3.4, 3.5,
   ['\u201c\u672a\u6765\u4e94\u5e74\u5185\u4f20\u7edf\u4e13\u4e1a\u4e0e\u8bfe\u7a0b\u53ef\u80fd\u6d88\u4ea1\uff0c\u6559\u80b2\u6838\u5fc3\u5e94\u8f6c\u5411\u57f9\u517b\u5b66\u751f\u5bf9AI\u751f\u6210\u5185\u5bb9\u7684\u5224\u65ad\u529b\u4e0e\u6279\u5224\u6027\u601d\u7ef4\u3002\u201d',
    '',
    '\u57fa\u4e8e\u4e09\u5927\u8d8b\u52bf\uff1a',
    '1. AI\u751f\u6210\u80fd\u529b\u6301\u7eed\u63d0\u5347',
    '2. \u77e5\u8bc6\u4f20\u6388\u4ef7\u503c\u4e0b\u964d',
    '3. \u5224\u65ad\u529b\u6210\u4e3a\u6838\u5fc3\u7ade\u4e89\u529b',
    '',
    '\u6838\u5fc3\u4e3b\u5f20\uff1a',
    '\u2022 \u4fe1\u5ff5\u5148\u884c',
    '\u2022 \u63a2\u7d22\u4e0d\u6b62',
    '\u2022 \u8303\u5f0f\u91cd\u6784'],
   fs=13, c=C['white'])

tb(s, 0.8, 6.0, 11, 0.8,
   '\u201c\u7ea0\u7ed3\u8bfe\u7a0b\u6709\u6ca1\u6709\u610f\u4e49\u5df2\u65e0\u5fc5\u8981\u3002\u201d\u2014\u2014 \u8bf8\u845b\u658c',
   fs=16, c=C['primary'], italic=True, a=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 7: Innovation Model
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, C['white'])
hdr(s, '\u4e94\u3001\u521b\u65b0\u5b9e\u8df5\u6a21\u5f0f\uff1a\u201c\u5b66\u751f\u4f5c\u4e1a\u5373\u8bfe\u7a0b\u8d44\u6e90\u201d')

tb(s, 0.8, 1.4, 4, 0.5, '\u64cd\u4f5c\u6d41\u7a0b', fs=20, c=C['accent'], b=True)

steps = [
    ('\u2460 \u77e5\u8bc6\u70b9\u62c6\u89e3', '\u5c06\u8bfe\u7a0b\u77e5\u8bc6\u70b9\u62c6\u89e3\u4e3a\u5b66\u751f\u4f5c\u4e1a\u4efb\u52a1\uff0c\u6bcf\u8282\u8bfe1-2\u4e2a\u77e5\u8bc6\u70b9'),
    ('\u2461 \u5b66\u751f\u5b9e\u8df5', '\u5b66\u751f\u4f7f\u7528AI\u5de5\u5177\uff08Manus\u3001\u667a\u80fd\u4f53\u7b49\uff09\u751f\u6210\u52a8\u753b\u3001\u8bb2\u89e3\u89c6\u9891\u3001\u4ee3\u7801\u7b49\u53ef\u89c6\u5316\u5185\u5bb9'),
    ('\u2462 \u6548\u679c\u7b5b\u9009', '\u6559\u5e08\u4ece\u5168\u73ed\u63d0\u4ea4\u7684AI\u751f\u6210\u5185\u5bb9\u4e2d\u6311\u9009\u6bcf\u77e5\u8bc6\u70b93\u4e2a\u6700\u4f18\u4f5c\u54c1'),
    ('\u2463 \u8d44\u6e90\u5e93\u5efa\u8bbe', '\u5f62\u6210\u53ef\u590d\u7528\u7684\u77e5\u8bc6\u70b9\u8d44\u6e90\u5e93\uff0c\u201c\u6bcf\u77e5\u8bc6\u70b9\u914d\u4e00\u4e2a\u4f18\u8d28AI\u751f\u6210\u5730\u5740\u201d')
]

sy = 2.0
for title, desc in steps:
    rect(s, 0.8, sy, 5.5, 0.9, C['light'], r=True)
    tb(s, 1.0, sy+0.05, 5.0, 0.35, title, fs=14, c=C['primary'], b=True)
    tb(s, 1.0, sy+0.4, 5.0, 0.45, desc, fs=11, c=C['text'])
    sy += 1.0

rect(s, 7, 1.3, 5.5, 3.8, C['light'], C['accent'], 2, True)
tb(s, 7.2, 1.4, 5, 0.5, '\u4ef7\u503c\u8bc4\u5224\u6743\u8f6c\u79fb', fs=18, c=C['primary'], b=True, a=PP_ALIGN.CENTER)
mt(s, 7.2, 2.0, 5.1, 3.0,
   ['\u5747\u503c\u5bfc\u5411 \u2192 \u5cf0\u503c\u53d1\u73b0',
    '  \u4e0d\u8981\u6c42\u6240\u6709\u4eba\u505a\u5bf9\uff0c\u800c\u662f\u53d1\u73b0\u4f18\u79c0\u4f5c\u54c1',
    '',
    '\u6559\u5e08\u4e3b\u89c2\u8bc4\u4ef7 \u2192 \u5b66\u751f\u4f5c\u54c1\u8d28\u91cf\u9a71\u52a8',
    '  \u4ee5\u4f5c\u54c1\u8d28\u91cf\u800c\u975e\u6559\u5e08\u504f\u597d\u4e3a\u6807\u51c6',
    '',
    '\u4e00\u6b21\u6027\u8bc4\u4ef7 \u2192 \u8fed\u4ee3\u4f18\u5316',
    '  \u4e0d\u8981\u6c42\u4e00\u6b21\u6027\u5b8c\u7f8e\uff0c\u901a\u8fc7\u6301\u7eed\u4f18\u5316\u63d0\u5347\u6548\u679c',
    '',
    '\u201c\u4e0d\u5728\u4e8eAI\u80fd\u5426\u5b8c\u6210\u4efb\u52a1\uff0c',
    '  \u800c\u5728\u4e8e\u82b1\u66f4\u591a\u65f6\u95f4\u540e\u662f\u5426\u80fd\u505a\u5f97\u66f4\u597d\u3002\u201d'],
   fs=13, c=C['text'])

rect(s, 0.8, 5.5, 11.5, 1.5, C['light'], C['secondary'], 1.5, True)
tb(s, 1, 5.55, 3, 0.4, 'SDP\u8bfe\u7a0b\u5b9e\u8df5\u6848\u4f8b', fs=16, c=C['secondary'], b=True)
mt(s, 1, 5.95, 11, 1.0,
   ['\u2022 \u77e5\u8bc6\u70b9\u62c6\u89e3\uff1a\u5c06SDP\u6838\u5fc3\u6982\u5ff5\uff08\u63a7\u5236\u5668\u3001\u6570\u636e\u5e73\u9762\u3001OpenFlow\u534f\u8bae\u7b49\uff09\u62c6\u89e3\u4e3a10\u4e2a\u4f5c\u4e1a\u4efb\u52a1',
    '\u2022 \u4ece100\u4efd\u4f5c\u4e1a\u4e2d\u7b5b\u9009\u51fa30\u4efd\u4f18\u8d28\u4f5c\u54c1\uff0c\u5f62\u6210SDP\u77e5\u8bc6\u70b9\u8d44\u6e90\u5e93\uff08\u52a8\u753b/\u8bb2\u89e3/\u4ee3\u7801\u793a\u4f8b\uff09',
    '\u2022 \u5173\u952e\u539f\u5219\uff1a\u4e0d\u6ee1\u8db3\u4e8e\u201c\u770b\u4e0a\u53bb\u6f02\u4eae\u201d\uff0c\u63a8\u52a8AI\u751f\u6210\u5185\u5bb9\u8fbe\u5230\u66f4\u9ad8\u5a92\u4f53\u5316\u6c34\u5e73\uff0c\u7b26\u5408\u4e13\u4e1a\u6807\u51c6'],
   fs=13, c=C['text'])

# ============================================================

# ============================================================
# SLIDE 7B: DingTalk AI Management (NEW - with screenshot)
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, C['white'])
hdr(s, '钉钉群AI原生课程管理实践')

# Left: mechanism description
tb(s, 0.6, 1.3, 5.8, 0.5, '管理机制', fs=20, c=C['primary'], b=True)
mt(s, 0.6, 1.85, 6, 4.2,
   ['▶ 基于钉钉群的AI助手协作模式',
    '',
    '• AI助手：“小龙虾”（基于OpenClaw+阿里云AI）',
    '• 群名称：“智能体小龙虾测试”',
    '• 参与者：课程团队全体成员 + AI助手',
    '',
    '▶ 核心能力',
    '',
    '• 会议纪要自动生成与归档',
    '• 教学材料自动生成与打包',
    '• 课程案例自动导出与归档',
    '• 24小时在线知识库答疑',
    '',
    '▶ 实践数据',
    '',
    '• 已生成7+个完整教学案例',
    '• 累计处理200+条消息',
    '• 生成156个文件（Word/PPT/ZIP）',
    '• 平均响应时间 < 30秒'],
   fs=13, c=C['text'])

# Right: screenshot
img_path = '/home/admin/.openclaw/workspace/dingtalk-group-full-screenshot.png'
if os.path.exists(img_path):
    s.shapes.add_picture(img_path, Inches(7.0), Inches(1.5), Inches(5.5), Inches(4.0))
    tb(s, 7.2, 5.6, 5.2, 0.4, '图：钉钉群AI协作场景——群成员@AI助手自动生成PPT汇报材料', fs=11, c=C['muted'], a=PP_ALIGN.CENTER)

# Bottom highlight
rect(s, 0.6, 6.1, 12, 0.9, C['light'], C['accent'], 2, True)
tb(s, 0.8, 6.15, 11.5, 0.4, '核心价值', fs=16, c=C['primary'], b=True)
mt(s, 0.8, 6.5, 11.5, 0.5,
   ['• 教师从繁琐的文档整理中解放，专注教学设计与质量把关 | • AI助手实现会议纪要、课程总结、案例导出等高效工作流'],
   fs=13, c=C['text'])


# ============================================================
# SLIDE 8: Five Reform Measures
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, C['white'])
hdr(s, '\u516d\u3001\u4e94\u5927\u6539\u9769\u4e3e\u63aa')

ini = [
    { 'n': '01', 't': '\u534f\u8bae\u673a\u5236\u53ef\u89c6\u5316',
      'd': '\u5229\u7528AI\u5de5\u5177\u751f\u6210\u52a8\u6001\u52a8\u753b\uff0c\u8986\u76d6STP\u3001\u8def\u7531\u67e5\u627e\u3001CHAP\u3001ACL\u3001NAT\u3001TCP\u4e09\u6b21\u63e1\u624b\u3001OpenFlow\u7b49\u6838\u5fc3\u534f\u8bae\u3002\u5b66\u751f\u4e3b\u5bfc\u751f\u6210\uff0c\u6559\u5e08\u7b5b\u9009\u6700\u4f18\u4f5c\u54c1\uff0c\u5f62\u6210\u534f\u8bae\u52a8\u753b\u5e93\u3002', 'c': C['secondary'] },
    { 'n': '02', 't': '\u8bfe\u7a0b\u79c1\u6709\u77e5\u8bc6\u5e93\u673a\u5668\u4eba',
      'd': '\u6574\u5408\u6559\u6750\u3001\u8bfe\u4ef6\u3001\u4e60\u9898\u3001FAQ\uff0c\u63d0\u4f9b24\u5c0f\u65f6\u4e2a\u6027\u5316\u8f85\u5bfc\u7b54\u7591\u3002\u6280\u672f\u5b9e\u73b0\uff1a\u57fa\u4e8e\u767e\u70bc\u77e5\u8bc6\u5e93\u00d7OpenClaw\u96c6\u6210\u65b9\u6848\uff0c\u5b9e\u73b0\u4e09\u7ea7\u8bb0\u5fc6\u67b6\u6784\uff08L1\u5de5\u4f5c\u8bb0\u5fc6\u3001L2\u77e5\u8bc6\u8bb0\u5fc6\u3001L3\u957f\u671f\u8bb0\u5fc6\uff09\u3002', 'c': C['teal'] },
    { 'n': '03', 't': '\u5168\u6d41\u7a0b\u667a\u6167\u6559\u5b66\u7ba1\u7406',
      'd': 'MOOC\u6155\u8bfe\u5802 + \u9489\u9489AI\u52a9\u7406\uff0c\u8bfe\u524d-\u8bfe\u4e2d-\u8bfe\u540e\u5168\u6d41\u7a0b\u8986\u76d6\uff0c\u6570\u636e\u9a71\u52a8\u7684\u6559\u5b66\u6539\u8fdb\u3002\u901a\u8fc7\u9489\u9489\u7fa4\u4e0eAI\u52a9\u624b\u534f\u4f5c\uff0c\u5b9e\u73b0\u4f1a\u8bae\u7eaa\u8981\u81ea\u52a8\u751f\u6210\u3001\u8bfe\u7a0b\u603b\u7ed3\u667a\u80fd\u6c47\u603b\u7b49\u9ad8\u6548\u5de5\u4f5c\u6d41\u3002', 'c': C['accent'] },
    { 'n': '04', 't': '\u667a\u80fd\u4f53\u8f85\u52a9\u5b9e\u9a8c\u65b0\u8303\u5f0f',
      'd': 'AI\u8f85\u52a9\u751f\u6210\u547d\u4ee4 + \u4eba\u5de5\u4e13\u4e1a\u6392\u9519 + \u771f\u5b9e\u8bbe\u5907\u90e8\u7f72\uff0c\u4fdd\u7559\u771f\u5b9e\u4ea4\u4e92\u903b\u8f91\uff0cAI\u8f85\u52a9\u6982\u5ff5\u7406\u89e3\u800c\u975e\u66ff\u4ee3\u5b9e\u9a8c\u64cd\u4f5c\u3002\u9a8c\u8bc1\u6027/\u8bbe\u8ba1\u6027/\u7efc\u5408\u6027=3:4:3\uff0c\u5b9e\u9a8c\u6548\u7387\u63d0\u534740%\u3002', 'c': C['success'] },
    { 'n': '05', 't': '\u667a\u80fd\u4f53\u9a71\u52a8\u7684\u8bfe\u7a0b\u95e8\u6237\u5efa\u8bbe',
      'd': '\u8bfe\u7a0b\u5b98\u7f51\uff0c\u96c6\u6210\u6559\u5b66\u8d44\u6e90+\u5b9e\u9a8c\u6848\u4f8b+\u4f18\u79c0\u4f5c\u54c1+\u667a\u80fd\u4e92\u52a8\u3002\u5a92\u4f53\u5316\u8868\u8fbe\uff0c\u8d85\u8d8a\u4f20\u7edfPPT\u548c\u6559\u6750\u7684\u89c6\u89c9\u5448\u73b0\u9ad8\u5ea6\u3002', 'c': C['purple'] }
]

iy = 1.3
for it in ini:
    rect(s, 0.6, iy, 0.7, 1.0, it['c'], r=True)
    tb(s, 0.65, iy+0.15, 0.6, 0.7, it['n'], fs=20, c=C['white'], b=True, a=PP_ALIGN.CENTER)
    tb(s, 1.5, iy+0.05, 4.5, 0.4, it['t'], fs=16, c=it['c'], b=True)
    tb(s, 1.5, iy+0.45, 11.0, 0.55, it['d'], fs=12, c=C['text'])
    iy += 1.1

# ============================================================
# SLIDE 9: Assessment
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, C['white'])
hdr(s, '\u4e03\u3001\u6559\u5b66\u6210\u6548\u8bc4\u4f30\u4e0e\u6311\u6218\u5bf9\u7b56')

tb(s, 0.8, 1.4, 5, 0.5, '\u8bc4\u4f30\u539f\u5219', fs=18, c=C['accent'], b=True)
rect(s, 0.8, 2.0, 5.5, 1.8, C['light'], C['accent'], 2, True)
mt(s, 0.9, 2.05, 5.2, 1.7,
   ['\u201c\u4e0d\u5728\u4e8eAI\u80fd\u5426\u5b8c\u6210\u4efb\u52a1\uff0c\u800c\u5728\u4e8e\'',
    '\u82b1\u66f4\u591a\u65f6\u95f4\u540e\u662f\u5426\u80fd\u505a\u5f97\u66f4\u597d\'\u3002\u201d',
    '',
    '\u2022 \u4e0d\u8981\u6c42\u4e00\u6b21\u6027\u5b8c\u7f8e',
    '\u2022 \u8fed\u4ee3\u601d\u7ef4\uff1a\u5bf9\u201c\u5783\u573e\u4f46\u53ef\u7528\u201d\u5185\u5bb9\u4fdd\u6301\u9ad8\u5bb9\u5fcd\u5ea6',
    '\u2022 \u6559\u80b2\u5b9e\u6548\u4f18\u5148'],
   fs=13, c=C['text'])

# Comparison table
tb(s, 7, 1.4, 5, 0.5, '\u8bc4\u4f30\u6307\u6807\u5bf9\u6bd4', fs=18, c=C['primary'], b=True)
rect(s, 7, 2.0, 5.5, 0.4, C['primary'])
tb(s, 7.1, 2.0, 2.7, 0.4, '\u6307\u6807', fs=12, c=C['white'], b=True)
tb(s, 9.8, 2.0, 1.3, 0.4, '\u4f20\u7edf\u6a21\u5f0f', fs=12, c=C['white'], b=True)
tb(s, 11.1, 2.0, 1.3, 0.4, 'AI\u539f\u751f', fs=12, c=C['white'], b=True)
rows = [
    ('\u4f5c\u4e1a\u5b8c\u6210\u5ea6', '100%', '\u4e0d\u8981\u6c42\u5168\u5458'),
    ('\u4f5c\u54c1\u8d28\u91cf', '\u6559\u5e08\u4e3b\u89c2', '\u4f5c\u54c1\u8d28\u91cf\u9a71\u52a8'),
    ('\u8d44\u6e90\u79ef\u7d2f', '\u6559\u5e08\u4e3b\u5bfc', '\u5b66\u751f\u4f5c\u4e1a\u5373\u8d44\u6e90'),
    ('\u5b66\u751f\u53c2\u4e0e\u5ea6', '\u88ab\u52a8\u63a5\u53d7', '\u4e3b\u52a8\u521b\u9020'),
    ('\u5224\u65ad\u529b\u57f9\u517b', '\u65e0', '\u6838\u5fc3\u76ee\u6807')
]
ry = 2.45
for i, (c1, c2, c3) in enumerate(rows):
    bg_c = C['white'] if i % 2 == 0 else C['light']
    rect(s, 7, ry, 5.5, 0.35, bg_c)
    tb(s, 7.1, ry, 2.7, 0.35, c1, fs=11, c=C['text'])
    tb(s, 9.8, ry, 1.3, 0.35, c2, fs=11, c=C['text'], a=PP_ALIGN.CENTER)
    tb(s, 11.1, ry, 1.3, 0.35, c3, fs=11, c=C['text'], a=PP_ALIGN.CENTER)
    ry += 0.38

# Challenges
tb(s, 0.8, 4.2, 5, 0.5, '\u6311\u6218\u4e0e\u5bf9\u7b56', fs=18, c=C['primary'], b=True)
challenges = [
    ('\u5185\u5bb9\u8d28\u91cf\u98ce\u9669', 'AI\u751f\u6210\u5185\u5bb9\u8d28\u91cf\u53c2\u5dee', '\u6559\u5e08\u4f5c\u4e3a\u8d28\u91cf\u628a\u5173\u8005\uff0c\u901a\u8fc7\u8fed\u4ee3\u4f18\u5316\u63d0\u5347\u8d28\u91cf'),
    ('\u5b9e\u9a8c\u8bfe\u9002\u914d\u6027', 'AI\u96be\u4ee5\u76f4\u63a5\u5d4c\u5165\u5b9e\u9a8c\u8bfe', '\u7406\u8bba\u90e8\u5206\u7ed3\u5408AI\u8f85\u52a9\uff0c\u64cd\u4f5c\u90e8\u5206\u4fdd\u7559\u771f\u5b9e\u4ea4\u4e92'),
    ('\u6559\u5e08\u89d2\u8272\u8f6c\u53d8', '\u4ece\u4f20\u6388\u8005\u5230\u5f15\u5bfc\u8005', '\u51cf\u5c11\u4f20\u7edfPPT\u7f16\u5199\uff0c\u805a\u7126\u7b5b\u9009\u548c\u5f15\u5bfc'),
    ('\u56fd\u9645\u5316\u9002\u914d', '\u82f1\u65b9\u6559\u5b66\u6750\u6599\u540c\u6b65\u590d\u6742', '\u5efa\u7acb\u4e2d\u82f1\u53cc\u8bed\u8d44\u6e90\u751f\u6210\u673a\u5236')
]
cy = 4.8
for title, chal, sol in challenges:
    rect(s, 0.8, cy, 11.5, 0.55, C['light'], r=True)
    tb(s, 1, cy+0.02, 2, 0.5, title, fs=13, c=C['primary'], b=True)
    tb(s, 3.5, cy+0.02, 3.5, 0.5, f'\u6311\u6218\uff1a{chal}', fs=11, c=C['muted'])
    tb(s, 7.5, cy+0.02, 4.5, 0.5, f'\u5bf9\u7b56\uff1a{sol}', fs=11, c=C['text'])
    cy += 0.6

# ============================================================
# SLIDE 10: Expected Outcomes
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, C['white'])
hdr(s, '\u516b\u3001\u9884\u671f\u6210\u6548\u4e0e\u7ed3\u8bba')

tb(s, 0.8, 1.4, 5, 0.5, '\u9884\u671f\u6210\u6548', fs=18, c=C['accent'], b=True)
outcomes = [
    '\u5f62\u6210\u534f\u8bae\u52a8\u753b\u5e93 + \u8bfe\u7a0b\u4e13\u5c5e\u79c1\u6709\u77e5\u8bc6\u5e93\u673a\u5668\u4eba',
    '\u66f4\u65b0\u51fa\u7248\u914d\u5957\u6570\u667a\u5316\u6559\u6750',
    '\u5b9e\u9a8c\u6548\u7387\u63d0\u534740%+\uff0c\u6bcf\u5b66\u671f\u53d7\u76ca200+\u5b66\u751f',
    '\u6253\u9020\u201c101\u8ba1\u5212\u201d\u672c\u5730\u5316\u201cAI+\u6559\u80b2\u201d\u5178\u578b\u6848\u4f8b',
    '\u53ef\u590d\u5236\u3001\u53ef\u63a8\u5e7f\uff1a\u5f62\u6210\u7cfb\u7edf\u5316\u3001\u9ad8\u8d28\u91cf\u7684\u8bfe\u7a0b\u8d44\u6e90\u5e93'
]
oy = 1.9
for o in outcomes:
    rect(s, 0.8, oy, 0.5, 0.5, C['accent'], r=True)
    tb(s, 1.5, oy-0.03, 10.5, 0.55, o, fs=15, c=C['text'])
    oy += 0.7

tb(s, 0.8, 5.5, 5, 0.5, '\u7814\u7a76\u7ed3\u8bba', fs=18, c=C['primary'], b=True)
mt(s, 0.8, 6.0, 11, 1.0,
   ['\u672c\u7814\u7a76\u63d0\u51fa\u5e76\u9a8c\u8bc1\u4e86AI\u539f\u751f\u6559\u5b66\u8303\u5f0f\u7684\u53ef\u884c\u6027\uff1a\u4ece\u201c\u77e5\u8bc6\u704c\u8f93\u201d\u8f6c\u5411\u201cAI\u5e94\u7528\u5b9e\u8df5\u201d\uff0c\u4ece\u201c\u77e5\u8bc6\u8bb0\u5fc6\u201d\u8f6c\u5411\u201c\u5224\u65ad\u529b\u57f9\u517b\u201d\uff0c\u6784\u5efa\u201c\u5b66\u751f\u4f5c\u4e1a\u5373\u8bfe\u7a0b\u8d44\u6e90\u201d\u7684\u521b\u65b0\u6a21\u5f0f\uff0c\u5f62\u6210\u4e94\u5927\u6539\u9769\u4e3e\u63aa\uff0c\u4e3a\u672a\u6765\u4e13\u4e1a\u6559\u80b2\u7684\u5168\u9762\u91cd\u6784\u63d0\u4f9b\u4e86\u53ef\u590d\u5236\u3001\u53ef\u63a8\u5e7f\u7684\u5b9e\u8df5\u8def\u5f84\u3002'],
   fs=14, c=C['text'])

# ============================================================
# SLIDE 11: Conclusion
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, C['white'])
hdr(s, '\u4e5d\u3001\u672a\u6765\u5c55\u671b')

rect(s, 0.8, 1.4, 11.5, 2.2, C['primary'], r=True)
tb(s, 1, 1.5, 11, 0.5, '\u5b9e\u8df5\u610f\u4e49', fs=18, c=C['accent'], b=True)
mt(s, 1, 2.0, 11, 1.5,
   ['1. \u53ef\u590d\u5236\u3001\u53ef\u63a8\u5e7f\uff1a\u5f62\u6210\u7cfb\u7edf\u5316\u3001\u9ad8\u8d28\u91cf\u7684\u8bfe\u7a0b\u8d44\u6e90\u5e93\uff0c\u4e3a\u5176\u4ed6\u8bfe\u7a0b\u5efa\u8bbe\u63d0\u4f9b\u53c2\u8003',
    '2. \u6559\u80b2\u5b9e\u6548\u4f18\u5148\uff1a\u6838\u5fc3\u5728\u4e8e\u6559\u80b2\u5b9e\u6548\u800c\u975e\u5f62\u5f0f\u7f8e\u89c2\uff0c\u907f\u514d\u201c\u770b\u4e0a\u53bb\u6f02\u4eae\u201d\u4f46\u7f3a\u4e4f\u6df1\u5ea6\u7684\u5185\u5bb9',
    '3. \u8fed\u4ee3\u601d\u7ef4\uff1a\u5bf9\u201c\u5783\u573e\u4f46\u53ef\u7528\u201d\u5185\u5bb9\u4fdd\u6301\u9ad8\u5bb9\u5fcd\u5ea6\uff0c\u901a\u8fc7\u7b5b\u9009\u4e0e\u53cd\u9988\u95ed\u73af\u9010\u6b65\u63d0\u7eaf\u4ef7\u503c'],
   fs=14, c=C['white'])

rect(s, 0.8, 4.0, 11.5, 2.2, C['light'], C['secondary'], 2, True)
tb(s, 1, 4.1, 11, 0.5, '\u4e94\u5e74\u540e\u5c55\u671b', fs=18, c=C['primary'], b=True)
mt(s, 1, 4.6, 11, 1.5,
   ['\u201c\u7ea0\u7ed3\u8bfe\u7a0b\u6709\u6ca1\u6709\u610f\u4e49\u5df2\u65e0\u5fc5\u8981\u3002\u201d',
    '',
    '\u4e94\u5e74\u540e\uff0c\u5f53AI\u751f\u6210\u80fd\u529b\u6301\u7eed\u63d0\u5347\uff0c\u4f20\u7edf\u4e13\u4e1a\u6559\u80b2\u5c06\u9762\u4e34\u5168\u9762\u91cd\u6784\u3002\u672c\u7814\u7a76\u7684\u610f\u4e49\u5728\u4e8e\uff1a',
    '\u2022 \u4fe1\u5ff5\u5148\u884c\uff1a\u6559\u5e08\u9700\u5148\u76f8\u4fe1AI\u7684\u6f5c\u529b\uff0c\u624d\u80fd\u8c03\u6574\u6559\u5b66\u7b56\u7565',
    '\u2022 \u63a2\u7d22\u4e0d\u6b62\uff1a\u662f\u5426\u80fd\u5b9e\u73b0\u4e0d\u91cd\u8981\uff0c\u91cd\u8981\u7684\u662f\u662f\u5426\u613f\u610f\u63a2\u7d22\u4e0e\u9a8c\u8bc1'],
   fs=14, c=C['text'])

tb(s, 0.8, 6.5, 11, 0.8,
   '\u201c\u4ea4\u4e00\u767e\u4efd\u4f5c\u4e1a\uff0c\u53ea\u8981\u6709\u4e24\u4efd\u8d85\u51fa\u60f3\u8c61\uff0c\u5c31\u662f\u6210\u529f\u3002\u201d \u2014\u2014 \u8bf8\u845b\u658c',
   fs=18, c=C['accent'], b=True, a=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 12: Thank You
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, C['primary'])
rect(s, 0, 0, 13.33, 0.12, C['accent'])
tb(s, 1, 1.5, 11.33, 1.0,
   '\u611f\u8c22\u8046\u542c\uff01',
   fs=44, c=C['white'], b=True, a=PP_ALIGN.CENTER)
tb(s, 1, 2.5, 11.33, 0.7,
   '\u6b22\u8fce\u6279\u8bc4\u6307\u6b63',
   fs=24, c=C['accent'], a=PP_ALIGN.CENTER)

rect(s, 2, 3.8, 9, 2.5, C['light'], r=True)
mt(s, 2.2, 3.9, 8.6, 2.3,
   ['\u57fa\u4e8eAI\u539f\u751f\u6559\u5b66\u8303\u5f0f\u7684\u8ba1\u7b97\u673a\u7f51\u7edc\u8bfe\u7a0b\u6539\u9769\u5b9e\u8df5',
    '',
    '\u57fa\u91d1\u9879\u76ee\uff1a\u6d59\u6c5f\u7701\u201cCS&AI 101\u8ba1\u5212\u201d\u8bfe\u7a0b\u5efa\u8bbe\u9879\u76ee\uff08\u6a21\u5757\u4e00\uff1a\u8ba1\u7b97\u673a\u7f51\u7edc\uff09',
    '',
    '\u4f5c\u8005\uff1a\u8bf8\u845b\u658c\u3001\u91d1\u84c9\u3001\u9ad8\u660e\u3001\u674e\u4f20\u714c\u3001\u848b\u732e',
    '\u901a\u4fe1\u4f5c\u8005\uff1a\u8bf8\u845b\u658c\uff0czhugubin@zjgsu.edu.cn'],
   fs=14, c=C['text'], a=PP_ALIGN.CENTER)

out = '/home/admin/.openclaw/workspace/101-paper-based-ppt.pptx'
prs.save(out)
print(f'PPT saved: {out}')
