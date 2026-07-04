#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
新信息技术赛道参赛项目 PPT 生成器
"""

from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Cm(33.867)
prs.slide_height = Cm(19.05)

COLORS = {
    'primary': (26, 35, 126),
    'secondary': (67, 97, 238),
    'accent': (0, 210, 255),
    'success': (0, 255, 127),
    'warning': (255, 193, 7),
    'light': (240, 245, 255),
    'white': (255, 255, 255),
    'dark': (20, 25, 50),
    'gray': (120, 130, 150),
}

def set_background(slide, color=COLORS['white']):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*color)

def add_header(slide, title):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(0), Cm(0), prs.slide_width, Cm(1.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*COLORS['primary'])
    shape.line.fill.background()
    
    textbox = slide.shapes.add_textbox(Cm(1.5), Cm(0.3), Cm(18), Cm(1))
    tf = textbox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(*COLORS['white'])
    p.font.name = 'Microsoft YaHei'

def add_footer(slide, page_num=""):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(2), prs.slide_height - Cm(1), prs.slide_width - Cm(4), Cm(0.05))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*COLORS['secondary'])
    shape.line.fill.background()
    
    textbox = slide.shapes.add_textbox(prs.slide_width - Cm(3), prs.slide_height - Cm(1.3), Cm(2), Cm(0.5))
    tf = textbox.text_frame
    p = tf.paragraphs[0]
    p.text = page_num
    p.font.size = Pt(11)
    p.font.color.rgb = RGBColor(*COLORS['gray'])
    p.font.name = 'Microsoft YaHei'
    p.alignment = PP_ALIGN.RIGHT

def add_title_slide(prs, title, subtitle, footer=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(0), Cm(0), prs.slide_width, Cm(5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*COLORS['primary'])
    shape.line.fill.background()
    
    textbox = slide.shapes.add_textbox(Cm(2), Cm(1.5), Cm(28), Cm(2.5))
    tf = textbox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(*COLORS['white'])
    p.font.name = 'Microsoft YaHei'
    
    textbox = slide.shapes.add_textbox(Cm(2), Cm(4.5), Cm(20), Cm(1.5))
    tf = textbox.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(22)
    p.font.color.rgb = RGBColor(*COLORS['accent'])
    p.font.name = 'Microsoft YaHei'
    
    if footer:
        textbox = slide.shapes.add_textbox(Cm(2), prs.slide_height - Cm(3), Cm(15), Cm(1))
        tf = textbox.text_frame
        p = tf.paragraphs[0]
        p.text = footer
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(*COLORS['gray'])
        p.font.name = 'Microsoft YaHei'
    
    return slide

def add_content_slide(prs, title, bullets, page_num=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_header(slide, title)
    
    y = Cm(2.5)
    for bullet in bullets:
        textbox = slide.shapes.add_textbox(Cm(2.5), y, Cm(29), Cm(0.7))
        tf = textbox.text_frame
        p = tf.paragraphs[0]
        p.text = bullet
        p.font.size = Pt(17)
        p.font.color.rgb = RGBColor(*COLORS['dark'])
        p.font.name = 'Microsoft YaHei'
        y += Cm(0.9)
    
    add_footer(slide, page_num)
    return slide

def add_diagram_slide(prs, title, layers, page_num=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_header(slide, title)
    
    y = Cm(3.5)
    for layer_title, layer_content, color in layers:
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(5), y, Cm(24), Cm(1.5))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*color)
        shape.fill.transparency = 0.7
        shape.line.color.rgb = RGBColor(*color)
        shape.line.width = Pt(2)
        
        textbox = slide.shapes.add_textbox(Cm(6), y + Cm(0.3), Cm(22), Cm(0.5))
        tf = textbox.text_frame
        p = tf.paragraphs[0]
        p.text = layer_title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = RGBColor(*COLORS['white'])
        p.font.name = 'Microsoft YaHei'
        
        textbox = slide.shapes.add_textbox(Cm(6), y + Cm(0.9), Cm(22), Cm(0.5))
        tf = textbox.text_frame
        p = tf.paragraphs[0]
        p.text = layer_content
        p.font.size = Pt(13)
        p.font.color.rgb = RGBColor(*COLORS['white'])
        p.font.name = 'Microsoft YaHei'
        
        y += Cm(2)
    
    add_footer(slide, page_num)
    return slide

def add_end_slide(prs, title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(0), Cm(0), prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*COLORS['primary'])
    shape.line.fill.background()
    
    textbox = slide.shapes.add_textbox(Cm(6), Cm(7), prs.slide_width - Cm(12), Cm(3))
    tf = textbox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = RGBColor(*COLORS['white'])
    p.font.name = 'Microsoft YaHei'
    p.alignment = PP_ALIGN.CENTER
    
    if subtitle:
        textbox = slide.shapes.add_textbox(Cm(6), Cm(10), prs.slide_width - Cm(12), Cm(2))
        tf = textbox.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(22)
        p.font.color.rgb = RGBColor(*COLORS['accent'])
        p.font.name = 'Microsoft YaHei'
        p.alignment = PP_ALIGN.CENTER
    
    return slide

# 生成幻灯片
add_title_slide(prs, "智学助手", "基于 AI 的个性化学习推荐系统", "新信息技术赛道参赛项目")

add_content_slide(prs, "项目目录", [
    "01 项目背景与痛点分析",
    "02 产品介绍与核心功能",
    "03 技术架构与创新亮点",
    "04 市场分析与竞争优势",
    "05 商业模式与盈利预测",
    "06 项目成果与应用成效",
    "07 团队介绍与未来规划",
], "01/12")

add_content_slide(prs, "项目背景", [
    "【学习痛点】",
    "  - 学习资源过载：海量资源难以选择",
    "  - 学习路径单一：无法满足个性化需求",
    "  - 学习效果难评估：缺乏精准诊断和反馈",
    "  - 学习动力不足：缺少针对性推荐",
    "",
    "【政策背景】",
    "  - 《教育信息化 2.0 行动计划》",
    "  - 《职业教育提质培优行动计划》",
    "  - 《新一代人工智能发展规划》",
], "02/12")

add_content_slide(prs, "产品介绍", [
    "【产品定位】",
    "智学助手是一款面向中职学生的 AI 个性化学习推荐系统",
    "",
    "【核心能力】",
    "  - 学情精准诊断",
    "  - 学习资源智能推荐",
    "  - 学习路径个性化规划",
    "  - 学习效果实时反馈",
    "",
    "【目标用户】",
    "  - 中职学生（15-18 岁）",
    "  - 高职学生（18-21 岁）",
    "  - 教师与学校",
], "03/12")

add_content_slide(prs, "核心功能", [
    "1. 学情诊断 - 分析学生知识掌握情况（知识图谱 + 诊断算法）",
    "2. 智能推荐 - 推荐适合的学习资源（协同过滤 + 内容推荐）",
    "3. 路径规划 - 生成个性化学习路径（强化学习算法）",
    "4. 效果追踪 - 实时监测学习效果（数据可视化 + 预警）",
    "5. 互动答疑 - 7x24 小时智能答疑（NLP+ 知识问答）",
], "04/12")

layers = [
    ("应用层", "Web/APP/小程序/管理后台", COLORS['secondary']),
    ("服务层", "推荐引擎/用户画像/学情分析/智能问答", COLORS['accent']),
    ("算法层", "协同过滤/知识图谱/NLP/机器学习", COLORS['success']),
    ("数据层", "MySQL/MongoDB/Redis/ES", COLORS['warning']),
    ("基础设施", "云服务器/容器/CDN/负载均衡", COLORS['primary']),
]
add_diagram_slide(prs, "技术架构", layers, "05/12")

add_content_slide(prs, "创新亮点", [
    "【技术创新】",
    "  - 融合多种推荐算法，推荐准确率 85%",
    "  - 构建职业教育专业知识图谱",
    "  - 基于强化学习的动态路径规划",
    "",
    "【模式创新】",
    "  - 诊断 - 推荐 - 学习 - 反馈闭环",
    "  - 教师+AI 双师协同模式",
    "",
    "【应用创新】",
    "  - 针对中职学生特点定制",
    "  - 对接职业技能标准",
], "06/12")

add_content_slide(prs, "市场分析", [
    "【目标市场】",
    "  - 全国中职学校：约 1 万所",
    "  - 中职在校生：约 1800 万人",
    "  - AI 教育市场：2025 年预计 1000 亿元",
    "",
    "【竞争优势】",
    "  - 智学助手：中职教育定位，推荐精度 85%",
    "  - 竞品 A: K12 定位，推荐精度 75%",
    "  - 竞品 B: 高等教育定位，推荐精度 80%",
], "07/12")

add_content_slide(prs, "商业模式", [
    "【收入来源】",
    "  - 基础服务：免费（吸引用户）",
    "  - 增值服务：会员订阅 19.9 元/月",
    "  - 学校采购：SaaS 服务 5 万元/校/年",
    "  - 数据服务：学情分析报告 2 万元/份",
    "",
    "【盈利预测】",
    "  - 第 1 年：用户 5 万，收入 50 万",
    "  - 第 2 年：用户 20 万，利润 50 万",
    "  - 第 3 年：用户 50 万，利润 250 万",
], "08/12")

add_content_slide(prs, "项目成果", [
    "【已实现功能】",
    "  - 用户注册与登录",
    "  - 学情诊断测试",
    "  - 智能推荐引擎 V1.0",
    "  - 学习路径规划",
    "  - 学习数据看板",
    "  - 智能问答机器人",
    "",
    "【技术成果】",
    "  - 软件著作权 2 项",
    "  - 专利申请 1 项",
    "  - 论文 1 篇",
], "09/12")

add_content_slide(prs, "应用成效", [
    "【试点学校】绍兴柯桥区高级技工学校",
    "",
    "【试点数据（3 个月）】",
    "  - 注册用户：1200 人",
    "  - 日活跃用户：450 人",
    "  - 平均使用时长：35 分钟/天",
    "  - 学习成效：平均成绩提高 15%",
    "  - 用户满意度：92%",
], "10/12")

add_content_slide(prs, "团队介绍", [
    "【核心成员】",
    "  - 项目负责人：XXX（计算机应用）",
    "  - 技术总监：XXX（软件工程）",
    "  - 产品总监：XXX（电子商务）",
    "  - 运营总监：XXX（市场营销）",
    "",
    "【指导教师】",
    "  - 何永胜（高级讲师，信息技术）",
    "  - XXX（副教授，人工智能）",
    "",
    "【顾问团队】教育专家、技术专家、行业专家",
], "11/12")

add_content_slide(prs, "未来规划", [
    "【短期目标（1 年）】",
    "  - 完善系统功能，拓展 10 所试点学校",
    "  - 用户突破 10 万，实现收支平衡",
    "",
    "【中期目标（3 年）】",
    "  - 覆盖浙江省 50% 中职学校",
    "  - 用户突破 50 万，年利润 500 万",
    "",
    "【长期目标（5 年）】",
    "  - 全国领先的职教 AI 平台",
    "  - 用户突破 500 万，科创板上市",
], "12/12")

add_end_slide(prs, "感谢观看", "敬请批评指正")

output_path = '/home/admin/.openclaw/workspace/新信息技术赛道参赛项目 PPT.pptx'
prs.save(output_path)
print(f"PPT 已生成：{output_path}")
print(f"共 {len(prs.slides)} 页")
