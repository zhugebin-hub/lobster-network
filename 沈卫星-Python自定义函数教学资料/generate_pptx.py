#!/usr/bin/env python3
"""生成《Python自定义函数》PPT课件（.pptx格式）"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import textwrap

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ===================== 配色方案 =====================
DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)       # 深蓝背景
ACCENT_BLUE = RGBColor(0x00, 0x96, 0xD6)    # 亮蓝
ACCENT_GREEN = RGBColor(0x00, 0xB4, 0xD8)   # 青绿
ACCENT_ORANGE = RGBColor(0xFF, 0x8C, 0x00)  # 橙色
ACCENT_RED = RGBColor(0xFF, 0x4D, 0x4D)     # 红色
ACCENT_YELLOW = RGBColor(0xFF, 0xD7, 0x00)  # 金色
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xE0, 0xE0, 0xE0)
CODE_BG = RGBColor(0x0D, 0x11, 0x17)        # 代码背景
CODE_GREEN = RGBColor(0x7C, 0xFC, 0x00)     # 代码绿色
CODE_YELLOW = RGBColor(0xFF, 0xD7, 0x00)    # 代码黄色
CODE_CYAN = RGBColor(0x00, 0xFF, 0xFF)      # 代码青色
CODE_PINK = RGBColor(0xFF, 0x69, 0xB4)      # 代码粉色
CODE_WHITE = RGBColor(0xE6, 0xED, 0xF3)     # 代码白色
BOX_BLUE = RGBColor(0x00, 0x7B, 0xFF)       # 信息框蓝色
BOX_GREEN = RGBColor(0x28, 0xA7, 0x45)      # 信息框绿色
BOX_ORANGE = RGBColor(0xFD, 0x7E, 0x14)     # 信息框橙色

# ===================== 工具函数 =====================

def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, left, top, width, height, text, font_size=18,
                font_color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                font_name="微软雅黑", line_spacing=1.2):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.line_spacing = Pt(font_size * line_spacing)
    return txBox

def add_code_block(slide, left, top, width, height, code_lines, font_size=14):
    """添加代码块"""
    # 背景
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = CODE_BG
    shape.line.fill.background()
    shape.shadow.inherit = False

    txBox = slide.shapes.add_textbox(left + Inches(0.15), top + Inches(0.1), width - Inches(0.3), height - Inches(0.2))
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, line in enumerate(code_lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.name = "Consolas"
        p.font.color.rgb = CODE_WHITE
        p.line_spacing = Pt(font_size * 1.5)
    return txBox

def add_bullet_list(slide, left, top, width, height, items, font_size=16,
                    font_color=WHITE, bullet_char="•", line_spacing=1.5):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"{bullet_char} {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.name = "微软雅黑"
        p.line_spacing = Pt(font_size * line_spacing)
    return txBox

def add_numbered_list(slide, left, top, width, height, items, font_size=16,
                      font_color=WHITE, line_spacing=1.5):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"{i+1}. {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.name = "微软雅黑"
        p.line_spacing = Pt(font_size * line_spacing)
    return txBox

def add_info_box(slide, left, top, width, height, title, content, color=BOX_BLUE, font_size=14):
    """添加信息框"""
    # 背景
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False

    # 标题
    add_textbox(slide, left + Inches(0.2), top + Inches(0.1), width - Inches(0.4),
                Inches(0.4), title, font_size=font_size+2, font_color=WHITE, bold=True)
    # 内容
    add_textbox(slide, left + Inches(0.2), top + Inches(0.5), width - Inches(0.4),
                height - Inches(0.6), content, font_size=font_size, font_color=WHITE)

def add_slide_number(slide, num, total=28):
    add_textbox(slide, Inches(12.3), Inches(7.0), Inches(0.8), Inches(0.4),
                f"{num}/{total}", font_size=10, font_color=LIGHT_GRAY,
                alignment=PP_ALIGN.RIGHT)

def add_section_header(slide, num, title, subtitle=""):
    set_slide_bg(slide, DARK_BG)
    # 顶部装饰线
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(0), Inches(0), Inches(13.333), Inches(0.08))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_BLUE
    line.line.fill.background()

    add_textbox(slide, Inches(1), Inches(2.5), Inches(11), Inches(1),
                f"第{num}页", font_size=16, font_color=ACCENT_BLUE, bold=True)
    add_textbox(slide, Inches(1), Inches(3.2), Inches(11), Inches(1.5),
                title, font_size=32, font_color=WHITE, bold=True)
    if subtitle:
        add_textbox(slide, Inches(1), Inches(4.8), Inches(11), Inches(0.8),
                    subtitle, font_size=18, font_color=LIGHT_GRAY)
    add_slide_number(slide, num)


# ===================== 第1页：封面 =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide, DARK_BG)

# 装饰线
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                               Inches(3), Inches(2.0), Inches(7.333), Inches(0.06))
line.fill.solid()
line.fill.fore_color.rgb = ACCENT_BLUE
line.line.fill.background()

add_textbox(slide, Inches(2), Inches(2.3), Inches(9), Inches(1),
            "Python 程序设计", font_size=44, font_color=WHITE, bold=True,
            alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(2), Inches(3.5), Inches(9), Inches(0.8),
            "第X章  自定义函数", font_size=32, font_color=ACCENT_BLUE, bold=True,
            alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(2), Inches(4.8), Inches(9), Inches(0.8),
            "—— 写好你的第一份「代码配方」——", font_size=20, font_color=LIGHT_GRAY,
            alignment=PP_ALIGN.CENTER)

# 底部信息
line2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                Inches(3), Inches(5.8), Inches(7.333), Inches(0.04))
line2.fill.solid()
line2.fill.fore_color.rgb = ACCENT_BLUE
line2.line.fill.background()

add_textbox(slide, Inches(2), Inches(6.0), Inches(9), Inches(0.5),
            "授课教师：沈卫星    课时：1课时（45分钟）", font_size=16, font_color=LIGHT_GRAY,
            alignment=PP_ALIGN.CENTER)

add_slide_number(slide, 1)


# ===================== 第2页：学习目标 =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, 2, "学习目标", "📌 本节课你要学会")

# 知识目标
add_info_box(slide, Inches(0.8), Inches(4.2), Inches(3.6), Inches(2.8),
             "📘 知识目标", "• 理解什么是函数\n• 掌握用 def 定义函数的语法\n• 理解函数参数和返回值", BOX_BLUE)

# 能力目标
add_info_box(slide, Inches(4.8), Inches(4.2), Inches(3.6), Inches(2.8),
             "🛠️ 能力目标", "• 能独立编写简单函数\n• 能调试常见函数错误\n• 能调用内置和自定义函数", BOX_GREEN)

# 素养目标
add_info_box(slide, Inches(8.8), Inches(4.2), Inches(3.6), Inches(2.8),
             "💡 素养目标", "• 培养模块化编程思维\n• 养成规范编码习惯\n• 提升逻辑分析能力", BOX_ORANGE)


# ===================== 第3页：导入 - 奶茶店 =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, 3, "导入：生活中的「函数」", "🧋 奶茶店点单 = 函数调用！")

# 左边：输入
add_info_box(slide, Inches(0.8), Inches(4.0), Inches(3.5), Inches(2.8),
             "📥 输入（参数）", "你告诉店员：\n• 杯型：大杯\n• 甜度：三分甜\n• 加料：珍珠", BOX_BLUE)

# 中间：处理
add_info_box(slide, Inches(4.8), Inches(4.0), Inches(3.5), Inches(2.8),
             "⚙️ 处理（函数体）", "店员根据配方：\n• 加茶底\n• 加糖\n• 加珍珠\n• 摇匀", BOX_GREEN)

# 右边：输出
add_info_box(slide, Inches(8.8), Inches(4.0), Inches(3.5), Inches(2.8),
             "📤 输出（返回值）", "端给你的奶茶 = 返回值\n\n💡 函数就是程序里的「配方」！", BOX_ORANGE)


# ===================== 第4页：函数的定义 =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, 4, "什么是函数？", "📖 函数的定义与三大要素")

# 定义
add_textbox(slide, Inches(0.8), Inches(4.0), Inches(11.5), Inches(1.0),
            "函数是一段 可重复使用 的代码块，它接收输入，进行处理，返回输出。",
            font_size=18, font_color=WHITE, bold=True)

# 三大要素
add_info_box(slide, Inches(0.8), Inches(5.2), Inches(3.4), Inches(1.8),
             "📥 输入", "参数（parameters）\n函数接收的数据", BOX_BLUE)

add_info_box(slide, Inches(4.8), Inches(5.2), Inches(3.4), Inches(1.8),
             "⚙️ 处理", "函数体（函数内部的代码）\n对数据进行处理", BOX_GREEN)

add_info_box(slide, Inches(8.8), Inches(5.2), Inches(3.4), Inches(1.8),
             "📤 输出", "返回值（return value）\n处理后的结果", BOX_ORANGE)


# ===================== 第5页：为什么需要函数 =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, 5, "为什么需要函数？", "🎯 函数的优势")

# 左边：没有函数
add_textbox(slide, Inches(0.8), Inches(4.0), Inches(5.5), Inches(0.4),
            "❌ 没有函数：重复代码满天飞", font_size=16, font_color=ACCENT_RED, bold=True)

add_code_block(slide, Inches(0.8), Inches(4.5), Inches(5.5), Inches(2.5), [
    'print("订单1：手机壳 × 2 = 31元")',
    'print("订单2：数据线 × 1 = 25元")',
    'print("订单3：充电宝 × 1 = 89元")',
    '# ... 100个订单要写100遍？！',
], font_size=12)

# 右边：有函数
add_textbox(slide, Inches(6.8), Inches(4.0), Inches(5.5), Inches(0.4),
            "✅ 有了函数：一次编写，反复使用", font_size=16, font_color=ACCENT_GREEN, bold=True)

add_code_block(slide, Inches(6.8), Inches(4.5), Inches(5.5), Inches(2.5), [
    "def show_order(name, qty, price):",
    "    total = qty * price",
    '    print(f"订单：{name} × {qty} = {total}元")',
    "",
    'show_order("手机壳", 2, 15.5)',
    'show_order("数据线", 1, 25.0)',
    'show_order("充电宝", 1, 89.0)',
    "# 写一次，调用无数次！",
], font_size=12)


# ===================== 第6页：函数的定义格式 =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, 6, "函数的定义格式", "📝 Python 中用 def 关键字定义函数")

add_code_block(slide, Inches(0.8), Inches(4.0), Inches(6.0), Inches(3.0), [
    "def 函数名(参数1, 参数2, ...):",
    '    """文档字符串：说明函数做什么"""',
    "    # 函数体（必须缩进4个空格！）",
    "    语句1",
    "    语句2",
    "    return 返回值  # 可选",
], font_size=14)

# 语法要点
add_textbox(slide, Inches(7.2), Inches(4.0), Inches(5.0), Inches(0.4),
            "📝 语法要点", font_size=18, font_color=ACCENT_YELLOW, bold=True)

add_numbered_list(slide, Inches(7.2), Inches(4.5), Inches(5.0), Inches(3.0), [
    "用 def 关键字开头",
    "函数名要见名知意（如 calc_price）",
    "括号后必须有冒号 :",
    "函数体必须缩进（4个空格）",
    "return 用于返回结果（可选）",
], font_size=15, font_color=WHITE)


# ===================== 第7页：第一个函数 =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, 7, "第一个函数：打招呼", "👋 编写你的第一个自定义函数")

add_code_block(slide, Inches(0.8), Inches(4.0), Inches(5.5), Inches(2.5), [
    "def greet():",
    '    """向用户打招呼"""',
    '    print("你好，欢迎来到Python世界！")',
    "",
    "# 定义函数不会自动执行！",
    "# 必须调用才会运行",
    "",
    "greet()   # ✅ 调用函数",
    "greet()   # ✅ 可以多次调用",
    "greet()   # ✅ 重复使用",
], font_size=13)

add_textbox(slide, Inches(6.8), Inches(4.0), Inches(5.5), Inches(0.4),
            "📤 运行结果：", font_size=16, font_color=ACCENT_GREEN, bold=True)

add_code_block(slide, Inches(6.8), Inches(4.5), Inches(5.5), Inches(1.5), [
    "你好，欢迎来到Python世界！",
    "你好，欢迎来到Python世界！",
    "你好，欢迎来到Python世界！",
], font_size=13)

add_info_box(slide, Inches(0.8), Inches(6.5), Inches(11.5), Inches(0.8),
             "⚠️ 注意", "定义函数 ≠ 执行函数！定义只是「写好配方」，调用才是「开始做菜」。", BOX_ORANGE)


# ===================== 第8页：函数参数 - 位置参数 =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, 8, "函数参数：位置参数", "📍 按顺序传递参数")

add_code_block(slide, Inches(0.8), Inches(4.0), Inches(5.5), Inches(2.5), [
    "def make_tea(name, sweetness):",
    '    """制作奶茶"""',
    '    print(f"制作一杯 {sweetness} 的 {name}")',
    "",
    "# 按顺序传递参数",
    'make_tea("珍珠奶茶", "三分甜")',
    'make_tea("绿茶", "不甜")',
    'make_tea("红豆奶茶", "五分甜")',
], font_size=13)

add_textbox(slide, Inches(6.8), Inches(4.0), Inches(5.5), Inches(0.4),
            "📤 运行结果：", font_size=16, font_color=ACCENT_GREEN, bold=True)

add_code_block(slide, Inches(6.8), Inches(4.5), Inches(5.5), Inches(1.5), [
    "制作一杯 三分甜 的 珍珠奶茶",
    "制作一杯 不甜 的 绿茶",
    "制作一杯 五分甜 的 红豆奶茶",
], font_size=13)

add_info_box(slide, Inches(0.8), Inches(6.5), Inches(11.5), Inches(0.8),
             "📝 位置参数要点", "• 按位置一一对应\n• 第一个值 → 第一个参数，第二个值 → 第二个参数\n• 顺序不能错！", BOX_BLUE)


# ===================== 第9页：函数参数 - 默认参数 =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, 9, "函数参数：默认参数", "🔄 给参数设置默认值")

add_code_block(slide, Inches(0.8), Inches(4.0), Inches(6.0), Inches(3.0), [
    "def make_tea(name, sweetness=\"五分甜\"):",
    '    """制作奶茶（默认五分甜）"""',
    '    print(f"制作一杯 {sweetness} 的 {name}")',
    "",
    "# 不传 sweetness → 使用默认值",
    'make_tea("奶茶")',
    '# 输出：制作一杯 五分甜 的 奶茶',
    "",
    "# 传入 sweetness → 覆盖默认值",
    'make_tea("奶茶", "三分甜")',
    '# 输出：制作一杯 三分甜 的 奶茶',
], font_size=12)

add_info_box(slide, Inches(7.2), Inches(4.0), Inches(5.0), Inches(3.0),
             "📝 默认参数要点", "• 定义时给参数赋默认值\n• 调用时可以不传该参数\n• 有默认值的参数要放在后面\n\n❌ 错误：\ndef f(a=1, b):  ← 默认参数不能在非默认参数前面！", BOX_ORANGE)


# ===================== 第10页：函数参数 - 关键字参数 =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, 10, "函数参数：关键字参数", "🔑 用参数名来传值，顺序随意！")

add_code_block(slide, Inches(0.8), Inches(4.0), Inches(6.0), Inches(3.0), [
    "def make_tea(name, sweetness=\"五分甜\", size=\"大杯\"):",
    '    """制作奶茶"""',
    '    print(f"制作一杯 {size} {sweetness} 的 {name}")',
    "",
    "# 用参数名来传值，顺序随意！",
    'make_tea(sweetness="不甜", name="绿茶", size="中杯")',
    'make_tea(name="奶茶")  # 其他用默认值',
    'make_tea("珍珠奶茶", size="大杯")  # 混合使用',
], font_size=12)

add_info_box(slide, Inches(7.2), Inches(4.0), Inches(5.0), Inches(2.5),
             "📝 关键字参数要点", "• 用 参数名=值 的方式传参\n• 顺序可以任意调换\n• 代码更清晰易读\n• 推荐：参数多的时候用关键字参数", BOX_BLUE)


# ===================== 第11页：电商场景 - 带参数的函数 =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, 11, "电商场景：带参数的函数", "🛒 浙江电商场景：计算商品总价")

add_code_block(slide, Inches(0.8), Inches(4.0), Inches(6.5), Inches(3.0), [
    "# 浙江电商场景：计算商品总价",
    "def calc_price(name, price, quantity=1):",
    '    """计算单个商品总价"""',
    "    total = price * quantity",
    '    print(f"📦 {name} × {quantity} = {total:.2f}元")',
    "    return total",
    "",
    "# 调用函数",
    'calc_price("手机壳", 15.5, 2)',
    '# 📦 手机壳 × 2 = 31.00元',
    "",
    'calc_price("数据线", 25.0)',
    '# 📦 数据线 × 1 = 25.00元',
], font_size=12)

add_info_box(slide, Inches(7.8), Inches(4.0), Inches(4.5), Inches(2.5),
             "💡 这个函数用了什么参数？", "• name：位置参数（必须传）\n• price：位置参数（必须传）\n• quantity：默认参数（不传=1）\n\n📌 电商实战场景，贴近工作实际！", BOX_GREEN)


# ===================== 第12页：函数返回值 - return =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, 12, "函数返回值：return 语句", "📤 return 的作用")

add_code_block(slide, Inches(0.8), Inches(4.0), Inches(5.5), Inches(2.5), [
    "def add(a, b):",
    '    """计算两个数的和"""',
    "    result = a + b",
    "    return result   # 返回结果",
    "",
    "# 接收返回值",
    "total = add(3, 5)",
    "print(total)   # 输出：8",
], font_size=13)

add_textbox(slide, Inches(6.8), Inches(4.0), Inches(5.5), Inches(0.4),
            "📝 return 的作用", font_size=18, font_color=ACCENT_YELLOW, bold=True)

add_numbered_list(slide, Inches(6.8), Inches(4.5), Inches(5.5), Inches(2.0), [
    "把结果「递出去」",
    "结束函数的执行",
    "返回值可以赋值给变量",
], font_size=16, font_color=WHITE)

add_info_box(slide, Inches(0.8), Inches(6.5), Inches(11.5), Inches(0.8),
             "🔑 关键理解", "return ≠ print  |  print → 在屏幕上显示（说给你听）  |  return → 把结果交给你（装进盒子）", BOX_BLUE)


# ===================== 第13页：return vs print =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, 13, "return vs print", "⚡ 理解 return 和 print 的区别")

add_code_block(slide, Inches(0.8), Inches(4.0), Inches(6.0), Inches(3.0), [
    "def with_return():",
    "    return 42",
    "",
    "def with_print():",
    "    print(42)",
    "",
    "# 调用两个函数",
    "a = with_return()",
    "b = with_print()",
    "",
    'print(f"a = {a}")   # a = 42',
    'print(f"b = {b}")   # b = None',
], font_size=13)

add_textbox(slide, Inches(7.2), Inches(4.0), Inches(5.0), Inches(0.4),
            "📤 运行结果：", font_size=16, font_color=ACCENT_GREEN, bold=True)

add_code_block(slide, Inches(7.2), Inches(4.5), Inches(5.0), Inches(1.5), [
    "42          ← with_print() 打印的",
    "a = 42      ← return 的值赋给了 a",
    "b = None    ← print 没有返回值，b 是 None",
], font_size=12)

add_info_box(slide, Inches(0.8), Inches(6.5), Inches(11.5), Inches(0.8),
             "⚠️ 记住", "return 的值可以存起来再用  |  print 的值「说完就没了」", BOX_ORANGE)


# ===================== 第14页：电商场景 - 订单计算 =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, 14, "电商场景：订单金额计算", "🛒 带返回值的函数实战")

add_code_block(slide, Inches(0.8), Inches(4.0), Inches(6.5), Inches(3.0), [
    "# 浙江电商场景：订单金额计算",
    "def calc_order(items):",
    '    """计算订单总金额"""',
    "    total = 0",
    "    for item in items:",
    "        total += item",
    "    return total",
    "",
    "# 调用函数",
    "cart = [15.5, 25.0, 89.0]  # 购物车商品",
    "total = calc_order(cart)",
    'print(f"订单总额：{total:.2f}元")',
    '# 输出：订单总额：129.50元',
], font_size=12)

add_info_box(slide, Inches(7.8), Inches(4.0), Inches(4.5), Inches(2.5),
             "💡 思考", "如果把 return 改成 print，\n还能用 total 继续计算折扣吗？\n\n→ 不能！因为 print 不返回值！", BOX_ORANGE)


# ===================== 第15页：折扣计算函数 =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, 15, "折扣计算函数", "💰 电商促销场景")

add_code_block(slide, Inches(0.8), Inches(4.0), Inches(6.5), Inches(3.0), [
    "def apply_discount(price, rate=0.9):",
    '    """应用折扣，返回折后价"""',
    "    final = price * rate",
    "    return final",
    "",
    "# 调用",
    "original = 129.50",
    "# 双十一 85 折",
    "final_price = apply_discount(original, 0.85)",
    'print(f"原价：{original:.2f}元")',
    'print(f"折后：{final_price:.2f}元")',
    'print(f"节省：{original - final_price:.2f}元")',
], font_size=12)

add_textbox(slide, Inches(7.8), Inches(4.0), Inches(4.5), Inches(0.4),
            "📤 运行结果：", font_size=16, font_color=ACCENT_GREEN, bold=True)

add_code_block(slide, Inches(7.8), Inches(4.5), Inches(4.5), Inches(1.5), [
    "原价：129.50元",
    "折后：110.08元",
    "节省：19.43元",
], font_size=13)


# ===================== 第16页：综合演示 =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, 16, "综合演示：电商促销计算器", "🎯 多个函数协同工作")

add_code_block(slide, Inches(0.5), Inches(3.8), Inches(12.3), Inches(3.5), [
    "# ============================ 电商促销计算器（完整示例）============================",
    "def calc_item_price(name, price, quantity=1):",
    '    """计算单个商品总价"""',
    "    total = price * quantity",
    '    print(f"📦 {name} × {quantity} = {total:.2f}元")',
    "    return total",
    "",
    "def apply_discount(total, rate=0.9):",
    '    """应用折扣"""',
    "    return total * rate",
    "",
    "def show_receipt(customer, items_total, final_price):",
    '    """显示结算单"""',
    '    print(f"\\n🧾 结算单")',
    '    print(f"👤 顾客：{customer}")',
    '    print(f"💰 商品总额：{items_total:.2f}元")',
    '    print(f"💰 实付金额：{final_price:.2f}元")',
    '    print("=" * 25)',
    "",
    "# 主程序",
    'customer = "小明同学"',
    'item1 = calc_item_price("手机壳", 15.5, 2)',
    'item2 = calc_item_price("数据线", 25.0, 1)',
    "items_total = item1 + item2",
    "final = apply_discount(items_total, 0.85)",
    "show_receipt(customer, items_total, final)",
], font_size=10)


# ===================== 第17页：综合演示 - 运行结果 =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, 17, "综合演示：运行结果", "🎯 模块化编程思想")

add_code_block(slide, Inches(0.8), Inches(4.0), Inches(5.5), Inches(1.5), [
    "📦 手机壳 × 2 = 31.00元",
    "📦 数据线 × 1 = 25.00元",
    "",
    "🧾 结算单",
    "👤 顾客：小明同学",
    "💰 商品总额：56.00元",
    "💰 实付金额：47.60元",
    "=========================",
], font_size=13)

add_textbox(slide, Inches(6.8), Inches(4.0), Inches(5.5), Inches(0.4),
            "🎯 这个程序用了几个函数？", font_size=16, font_color=ACCENT_YELLOW, bold=True)

add_numbered_list(slide, Inches(6.8), Inches(4.5), Inches(5.5), Inches(2.0), [
    "calc_item_price() → 计算商品价",
    "apply_discount()   → 计算折扣",
    "show_receipt()     → 显示结算单",
], font_size=15, font_color=WHITE)

add_textbox(slide, Inches(6.8), Inches(5.8), Inches(5.5), Inches(0.4),
            "🎯 体现了什么编程思想？", font_size=16, font_color=ACCENT_YELLOW, bold=True)

add_numbered_list(slide, Inches(6.8), Inches(6.2), Inches(5.5), Inches(1.2), [
    "模块化：每个函数做一件事",
    "可复用：函数可以被多次调用",
    "易维护：修改某个功能只需改对应函数",
], font_size=14, font_color=WHITE)


# ===================== 第18页：常见错误（一）语法错误 =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, 18, "常见错误（一）：语法错误", "🐛 初学者最常犯的错误")

add_code_block(slide, Inches(0.8), Inches(4.0), Inches(6.0), Inches(3.0), [
    "# ❌ 错误1：忘记冒号",
    "def greet()",
    '    print("Hello!")',
    "",
    "# ✅ 正确",
    "def greet():",
    '    print("Hello!")',
    "",
    "",
    "# ❌ 错误2：缩进不对",
    "def greet():",
    'print("Hello!")    # ← 没有缩进！',
    "",
    "# ✅ 正确",
    "def greet():",
    '    print("Hello!")  # ← 缩进4个空格',
], font_size=12)

add_info_box(slide, Inches(7.2), Inches(4.0), Inches(5.0), Inches(2.5),
             "🔍 Python 3.12+ 的错误提示更友好了！", "旧版本：\n  SyntaxError: invalid syntax\n\nPython 3.12+：\n  SyntaxError: expected \':\'\n  def greet()\n             ^\n💡 直接告诉你缺了什么！", BOX_GREEN)


# ===================== 第19页：常见错误（二）调用错误 =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, 19, "常见错误（二）：调用错误", "🐛 函数调用时的常见错误")

add_code_block(slide, Inches(0.8), Inches(4.0), Inches(6.0), Inches(3.0), [
    "# ❌ 错误3：函数名拼写错误",
    "def say_hello():",
    '    print("Hello!")',
    "",
    'say_helo()   # ← 拼错了！',
    "# 🔧 解决：仔细检查函数名",
    "",
    "",
    "# ❌ 错误4：参数数量不匹配",
    "def make_tea(name, sweetness):",
    '    print(f"{sweetness}的{name}")',
    "",
    'make_tea("奶茶")   # ← 少传了一个参数！',
    "# 🔧 解决：检查参数数量，或使用默认参数",
], font_size=12)

add_info_box(slide, Inches(7.2), Inches(4.0), Inches(5.0), Inches(2.5),
             "🐛 调试技巧", "① 看错误信息（Python会告诉你哪一行错了）\n② 逐行检查（从报错行往上找）\n③ 用 print 查看中间结果\n④ 把大问题拆成小步骤测试", BOX_BLUE)


# ===================== 第20页：常见错误（三）逻辑错误 =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, 20, "常见错误（三）：逻辑错误", "🐛 代码能运行但结果不对")

add_code_block(slide, Inches(0.8), Inches(4.0), Inches(6.0), Inches(3.0), [
    "# ❌ 错误5：忘记 return",
    "def add(a, b):",
    "    result = a + b",
    "    # ← 忘了写 return！",
    "",
    "total = add(3, 5)",
    "print(total)   # 输出：None（不是8！）",
    "",
    "# ✅ 正确",
    "def add(a, b):",
    "    result = a + b",
    "    return result",
    "",
    "",
    "# ❌ 错误6：return 放错位置",
    "def max_num(a, b):",
    "    if a > b:",
    "        return a",
    "    return b   # ← 这个位置才对",
], font_size=12)

add_info_box(slide, Inches(7.2), Inches(4.0), Inches(5.0), Inches(2.5),
             "💡 逻辑错误最难发现！", "代码能运行，但结果不对。\n\n排查方法：\n• 用 print 打印中间变量\n• 逐步跟踪代码执行\n• 用简单的测试数据验证\n• 检查 return 的位置", BOX_ORANGE)


# ===================== 第21页：Python 3.12+ 新特性 =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, 21, "Python 3.12+ 新特性", "🆕 函数相关的新特性")

add_numbered_list(slide, Inches(0.8), Inches(4.0), Inches(5.5), Inches(3.0), [
    "更好的错误提示",
    "   · 语法错误会精确指出位置",
    "   · 用 ^ 标记出错的地方",
    "   · 对初学者更友好！",
    "",
    "类型提示改进（Type Hints）",
    "   def calc(name: str, price: float) -> float:",
    "       return price",
    "",
    "f-string 中的 = 符号（调试利器）",
    "   x = 42",
    '   print(f"{x=}")   # 输出：x=42',
], font_size=14, font_color=WHITE)

add_info_box(slide, Inches(6.8), Inches(4.0), Inches(5.5), Inches(2.5),
             "💡 这些特性让写函数更简单！", "Python 3.12+ 让代码更清晰、\n调试更容易、错误提示更友好。\n\n推荐升级到最新版本！", BOX_GREEN)


# ===================== 第22页：实践任务 =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, 22, "上机实践任务", "🎯 动手写代码！时间：10分钟")

# 基础任务
add_info_box(slide, Inches(0.8), Inches(4.0), Inches(3.6), Inches(3.0),
             "📌 基础任务（必做）", "定义一个 greet(name) 函数，\n接收用户名，打印欢迎信息。\n\n示例：\ngreet(\"小明\")\n→ 欢迎小明！", BOX_BLUE)

# 进阶任务
add_info_box(slide, Inches(4.8), Inches(4.0), Inches(3.6), Inches(3.0),
             "📌 进阶任务（必做）", "定义一个 calc_discount(price, rate) 函数，\n接收原价和折扣率，返回折后价。\n\n示例：\ncalc_discount(100, 0.85)\n→ 85.0", BOX_GREEN)

# 挑战任务
add_info_box(slide, Inches(8.8), Inches(4.0), Inches(3.6), Inches(3.0),
             "📌 挑战任务（选做）", "定义一个 generate_receipt(customer, items, rate) 函数，\n接收顾客名、商品列表和折扣率，\n生成完整的结算单。", BOX_ORANGE)

add_info_box(slide, Inches(0.8), Inches(6.5), Inches(11.5), Inches(0.8),
             "🤝 提示", "可以同桌讨论，但代码要自己写！遇到问题先自己调试，实在解决不了再举手。", BOX_BLUE)


# ===================== 第23页：课堂小结 =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, 23, "课堂小结", "📋 今天学了什么？")

add_numbered_list(slide, Inches(0.8), Inches(4.0), Inches(11.5), Inches(3.5), [
    "函数是什么？  →  可重复使用的代码块 = 代码配方",
    "怎么定义函数？  →  def 函数名(参数): 函数体  return 返回值",
    "三种参数  →  位置参数（按顺序）、默认参数（有默认值）、关键字参数（按名字）",
    "返回值  →  return 返回结果，无 return 返回 None",
    "常见错误  →  冒号、缩进、参数匹配、return 位置",
], font_size=18, font_color=WHITE)


# ===================== 第24页：课后作业 =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, 24, "课后作业", "📝 巩固练习")

add_numbered_list(slide, Inches(0.8), Inches(4.0), Inches(11.5), Inches(3.0), [
    "基础题：编写一个函数 is_even(n)，判断一个数是否为偶数，返回 True/False。",
    "",
    "提高题：编写一个函数 calc_shipping(weight, distance)，根据重量和距离计算运费：",
    "   · 1kg以内 8元  · 超过1kg，每kg加3元  · 超过100km，额外加5元",
    "",
    "拓展题：编写一个简易计算器，包含四个函数：",
    "   add(a,b)、subtract(a,b)、multiply(a,b)、divide(a,b)",
], font_size=16, font_color=WHITE)

add_info_box(slide, Inches(0.8), Inches(6.5), Inches(11.5), Inches(0.8),
             "💡 下节课预告", "函数的嵌套调用与作用域  |  局部变量 vs 全局变量  |  递归函数", BOX_BLUE)


# ===================== 第25页：参数对比表 =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, 25, "三种参数对比", "📊 位置参数 vs 默认参数 vs 关键字参数")

add_code_block(slide, Inches(0.8), Inches(4.0), Inches(11.5), Inches(3.0), [
    "┌──────────────┬──────────────┬────────────────────────────────┐",
    "│   参数类型   │   写法       │           示例                 │",
    "├──────────────┼──────────────┼────────────────────────────────┤",
    "│  位置参数    │  按顺序传    │  f(\"奶茶\",\"三分甜\")          │",
    "│  默认参数    │  有默认值    │  f(\"奶茶\")                    │",
    "│  关键字参数  │  按名字传    │  f(sweetness=\"不甜\",         │",
    "│              │              │    name=\"绿茶\")               │",
    "└──────────────┴──────────────┴────────────────────────────────┘",
    "",
    "📌 参数定义顺序规则：",
    "   普通参数 → 默认参数",
    "   def f(a, b, c=10):  ✅",
    "   def f(a=10, b):     ❌  错误！",
], font_size=13)


# ===================== 第26页：无返回值函数 =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, 26, "无返回值的函数", "📤 没有 return 的函数")

add_code_block(slide, Inches(0.8), Inches(4.0), Inches(5.5), Inches(2.5), [
    "def show_message(name):",
    '    """只显示信息，不返回结果"""',
    '    print(f"欢迎 {name} 来到电商直播间！")',
    "",
    'result = show_message("小明")',
    'print(f"返回值是：{result}")',
], font_size=13)

add_textbox(slide, Inches(6.8), Inches(4.0), Inches(5.5), Inches(0.4),
            "📤 运行结果：", font_size=16, font_color=ACCENT_GREEN, bold=True)

add_code_block(slide, Inches(6.8), Inches(4.5), Inches(5.5), Inches(1.5), [
    "欢迎 小明 来到电商直播间！",
    "返回值是：None",
], font_size=13)

add_info_box(slide, Inches(0.8), Inches(6.5), Inches(11.5), Inches(0.8),
             "📝 说明", "没有 return 语句的函数 → 默认返回 None  |  None 表示「什么都没有」  |  适合只需要执行操作、不需要返回结果的场景", BOX_BLUE)


# ===================== 第27页：调试技巧 =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, 27, "调试技巧", "🐛 如何快速找到并修复错误")

add_numbered_list(slide, Inches(0.8), Inches(4.0), Inches(5.5), Inches(3.0), [
    "看错误信息",
    "   Python 会告诉你哪一行错了",
    "   仔细阅读错误类型和提示",
    "",
    "逐行检查",
    "   从报错行往上找原因",
    "   检查语法、拼写、缩进",
    "",
    "用 print 调试",
    "   在关键位置打印变量值",
    "   查看程序执行到哪一步",
    "",
    "拆分测试",
    "   把大问题拆成小步骤",
    "   每一步都测试验证",
], font_size=14, font_color=WHITE)


# ===================== 第28页：结束页 =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

# 装饰线
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                               Inches(3), Inches(2.5), Inches(7.333), Inches(0.06))
line.fill.solid()
line.fill.fore_color.rgb = ACCENT_BLUE
line.line.fill.background()

add_textbox(slide, Inches(2), Inches(3.0), Inches(9), Inches(1),
            "谢谢聆听！", font_size=44, font_color=WHITE, bold=True,
            alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(2), Inches(4.2), Inches(9), Inches(0.8),
            "Python 程序设计 · 自定义函数", font_size=24, font_color=ACCENT_BLUE,
            alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(2), Inches(5.2), Inches(9), Inches(0.6),
            "授课教师：沈卫星", font_size=20, font_color=LIGHT_GRAY,
            alignment=PP_ALIGN.CENTER)

# 底部装饰线
line2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                Inches(3), Inches(6.0), Inches(7.333), Inches(0.04))
line2.fill.solid()
line2.fill.fore_color.rgb = ACCENT_BLUE
line2.line.fill.background()

add_slide_number(slide, 28)


# ===================== 保存 =====================
output_path = "/home/admin/.openclaw/workspace/沈卫星-Python自定义函数教学资料/Python自定义函数.pptx"
prs.save(output_path)
print(f"✅ PPT 已保存: {output_path}")
print(f"📊 共 {len(prs.slides)} 页")
