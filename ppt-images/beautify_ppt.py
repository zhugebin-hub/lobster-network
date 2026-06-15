#!/usr/bin/env python3
"""美化 PPT - 添加背景图片和装饰元素"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

INPUT_PPTX = "/home/admin/.openclaw/media/inbound/71a9c9bb-3e0a-4230-a874-65fa73908788.pptx"
OUTPUT_PPTX = "/home/admin/.openclaw/workspace/中美贸易战的成因、影响与对策分析_美化版.pptx"
IMAGE_DIR = "/home/admin/.openclaw/workspace/ppt-images"

# 图片映射
IMAGE_MAP = {
    1: "chess.jpg",  # 封面 - 战略博弈
    2: "chess.jpg",  # 核心结论
    5: "trade_port.jpg",  # 经济因素
    6: "us_capitol.jpg",  # 政治因素
    7: "globe_strategy.jpg",  # 战略因素
    9: "economy.jpg",  # 宏观影响
    10: "trade_data.jpg",  # 贸易影响
    11: "chip.jpg",  # 产业影响
    12: "global_supply.jpg",  # 全球影响
    14: "china_develop.jpg",  # 中国应对
    15: "economy.jpg",  # 美国应对
    16: "international.jpg",  # 国际社会
    17: "handshake.jpg",  # 结论展望
}

def add_image_with_overlay(slide, image_path, left, top, width, height, overlay_alpha=0.4):
    """添加图片并叠加半透明层"""
    # 添加图片
    slide.shapes.add_picture(
        image_path,
        left=left,
        top=top,
        width=width,
        height=height
    )
    
    # 添加半透明覆盖层
    overlay = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left=left,
        top=top,
        width=width,
        height=height
    )
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = RGBColor(0, 0, 0)
    overlay.line.fill.background()
    
    return overlay

def set_full_background(slide, image_path):
    """设置全屏背景"""
    slide.background.fill.background()
    slide.shapes.add_picture(
        image_path,
        left=Inches(0),
        top=Inches(0),
        width=Inches(13.33),
        height=Inches(7.50)
    )
    
    # 添加半透明黑色覆盖层
    overlay = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left=Inches(0),
        top=Inches(0),
        width=Inches(13.33),
        height=Inches(7.50)
    )
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = RGBColor(0, 0, 0)
    overlay.line.fill.background()

def beautify_slide(slide, slide_num, image_name):
    """美化单个幻灯片"""
    if not image_name:
        return
    
    image_path = os.path.join(IMAGE_DIR, image_name)
    if not os.path.exists(image_path):
        print(f"  ⚠️  图片不存在: {image_path}")
        return
    
    print(f"  🖼️  使用图片: {image_name}")
    
    # 封面页 - 全屏背景
    if slide_num == 1:
        set_full_background(slide, image_path)
        print(f"  ✅ 封面页全屏背景")
        return
    
    # 内容页 - 右侧图片区域
    # 检查是否有图表
    has_chart = False
    content_height = 0
    for shape in slide.shapes:
        if hasattr(shape, 'has_chart') and shape.has_chart:
            has_chart = True
        # 计算内容高度
        if hasattr(shape, 'top') and hasattr(shape, 'height'):
            bottom = shape.top + shape.height
            if bottom > content_height:
                content_height = bottom
    
    # 在右侧添加装饰图片
    if not has_chart:
        try:
            add_image_with_overlay(
                slide, image_path,
                left=Inches(8.0),
                top=Inches(0.5),
                width=Inches(5.0),
                height=Inches(3.5)
            )
            print(f"  ✅ 右侧图片区域")
        except Exception as e:
            print(f"  ⚠️  添加图片失败: {e}")
    else:
        # 有图表的页面，在底部添加
        try:
            add_image_with_overlay(
                slide, image_path,
                left=Inches(0.5),
                top=Inches(5.2),
                width=Inches(12.0),
                height=Inches(2.0)
            )
            print(f"  ✅ 底部图片区域")
        except Exception as e:
            print(f"  ⚠️  添加图片失败: {e}")

def main():
    print("🎨 开始美化 PPT...\n")
    
    prs = Presentation(INPUT_PPTX)
    print(f"📄 幻灯片总数: {len(prs.slides)}\n")
    
    for i, slide in enumerate(prs.slides, 1):
        print(f"📝 第 {i} 页...")
        image_name = IMAGE_MAP.get(i)
        beautify_slide(slide, i, image_name)
    
    # 保存
    prs.save(OUTPUT_PPTX)
    size_mb = os.path.getsize(OUTPUT_PPTX) / 1024 / 1024
    print(f"\n✅ 美化完成！")
    print(f"📁 保存路径: {OUTPUT_PPTX}")
    print(f"📦 文件大小: {size_mb:.1f} MB")

if __name__ == "__main__":
    main()
