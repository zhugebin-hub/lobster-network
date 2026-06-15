#!/usr/bin/env python3
"""生成「民国文人的宗教观」PPT"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import copy

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ============ 配色方案 ============
BG_DARK = RGBColor(0x1B, 0x2A, 0x4A)    # 深蓝背景
BG_ACCENT = RGBColor(0x2C, 0x3E, 0x6B)   # 中蓝
TEXT_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_LIGHT = RGBColor(0xCC, 0xD5, 0xE5)
GOLD = RGBColor(0xE8, 0xC5, 0x47)         # 金色强调
ACCENT_RED = RGBColor(0xC0, 0x39, 0x2B)
ACCENT_TEAL = RGBColor(0x16, 0xA0, 0x85)
ACCENT_ORANGE = RGBColor(0xE6, 0x7E, 0x22)
ACCENT_PURPLE = RGBColor(0x8E, 0x44, 0xAD)
ACCENT_GREEN = RGBColor(0x27, 0xAE, 0x60)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT = RGBColor(0x2C, 0x3E, 0x50)
LIGHT_BG = RGBColor(0xF5, 0xF6, 0xFA)
CARD_BG = RGBColor(0xFF, 0xFF, 0xFF)

def set_slide_bg(slide, color):
    """设置纯色背景"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, fill_color, line_color=None):
    """添加矩形色块"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_rounded_rect(slide, left, top, width, height, fill_color, line_color=None):
    """添加圆角矩形"""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18,
                 font_color=TEXT_WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name="微软雅黑", line_spacing=None):
    """添加文本框"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    if line_spacing:
        p.space_after = Pt(line_spacing)
    return txBox

def add_bullet_list(slide, left, top, width, height, items, font_size=16,
                    font_color=TEXT_LIGHT, bullet_color=GOLD, font_name="微软雅黑",
                    spacing_after=8):
    """添加带金色圆点的列表"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.name = font_name
        p.space_after = Pt(spacing_after)
        # 设置段落级别的project bullet
        pPr = p._pPr
        if pPr is None:
            from pptx.oxml.ns import qn
            pPr = p._p.get_or_add_pPr()
        from pptx.oxml.ns import qn
        buChar = pPr.makeelement(qn('a:buChar'), {'char': '●'})
        # remove existing bullets
        for child in list(pPr):
            if child.tag.endswith('buChar') or child.tag.endswith('buNone'):
                pPr.remove(child)
        pPr.append(buChar)
        # color the bullet
        buClr = pPr.makeelement(qn('a:buClr'), {})
        srgbClr = buClr.makeelement(qn('a:srgbClr'), {'val': '%02X%02X%02X' % (bullet_color[0], bullet_color[1], bullet_color[2])})
        buClr.append(srgbClr)
        pPr.append(buClr)
        indent = pPr.makeelement(qn('a:indent'), {'indent': str(Pt(0))})
    return txBox

def add_accent_bar(slide, left, top, height, color):
    """添加左侧竖条装饰"""
    add_rect(slide, left, top, Inches(0.08), height, color)

def add_number_circle(slide, left, top, size, number, bg_color, text_color=WHITE):
    """添加数字圆圈"""
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = str(number)
    p.font.size = Pt(24)
    p.font.color.rgb = text_color
    p.font.bold = True
    p.font.name = "微软雅黑"
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(0)
    tf.paragraphs[0].space_after = Pt(0)
    return shape

# ==================== 幻灯片 1：封面 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide, BG_DARK)

# 顶部装饰线
add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), GOLD)
# 底部装饰线
add_rect(slide, Inches(0), Inches(7.44), Inches(13.333), Inches(0.06), GOLD)
# 左侧装饰竖条
add_rect(slide, Inches(0.8), Inches(2.0), Inches(0.06), Inches(3.5), GOLD)

# 主标题
add_text_box(slide, Inches(1.3), Inches(2.2), Inches(10), Inches(1.5),
             "民国文人的宗教观", font_size=48, font_color=TEXT_WHITE, bold=True,
             font_name="微软雅黑")

# 副标题
add_text_box(slide, Inches(1.3), Inches(3.8), Inches(10), Inches(0.8),
             "传统与现代之间的思想图景", font_size=24, font_color=GOLD,
             font_name="微软雅黑")

# 分隔线
add_rect(slide, Inches(1.3), Inches(4.8), Inches(4), Inches(0.03), GOLD)

# 作者/日期
add_text_box(slide, Inches(1.3), Inches(5.1), Inches(6), Inches(0.5),
             "戴建华  ｜  2026年6月", font_size=16, font_color=TEXT_LIGHT,
             font_name="微软雅黑")

# 右下角小装饰
add_text_box(slide, Inches(10), Inches(6.2), Inches(3), Inches(0.5),
             "思想史 · 宗教学 · 文化研究", font_size=13, font_color=RGBColor(0x66, 0x77, 0x99),
             font_name="微软雅黑", alignment=PP_ALIGN.RIGHT)

# ==================== 幻灯片 2：目录 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), GOLD)

# 标题
add_text_box(slide, Inches(1), Inches(0.5), Inches(5), Inches(0.8),
             "目  录", font_size=36, font_color=TEXT_WHITE, bold=True)
add_rect(slide, Inches(1), Inches(1.3), Inches(2), Inches(0.03), GOLD)

chapters = [
    ("壹", "引言：时代背景与问题意识", BG_ACCENT),
    ("贰", "儒家为体：以儒学为信仰", RGBColor(0xC0, 0x39, 0x2B)),
    ("叁", "佛教徒：皈依或深度修习", RGBColor(0x8E, 0x44, 0xAD)),
    ("肆", "基督徒：信仰或文化认同", RGBColor(0x2C, 0x3E, 0x6B)),
    ("伍", "道家/道教：隐逸或养生", RGBColor(0x16, 0xA0, 0x85)),
    ("陆", "无神论者 / 反宗教者", RGBColor(0xE6, 0x7E, 0x22)),
    ("柒", "神秘主义 / 泛神论者", RGBColor(0xE7, 0x4C, 0x3C)),
    ("捌", "学术化研究：宗教作为学问", RGBColor(0x27, 0xAE, 0x60)),
    ("玖", "总结：共性与当代启示", GOLD),
]

for i, (num, title, color) in enumerate(chapters):
    y = Inches(1.7) + Inches(i * 0.6)
    # 数字
    add_text_box(slide, Inches(1.2), y, Inches(0.6), Inches(0.5),
                 num, font_size=20, font_color=color, bold=True)
    # 标题
    add_text_box(slide, Inches(1.9), y, Inches(8), Inches(0.5),
                 title, font_size=18, font_color=TEXT_LIGHT)

# ==================== 幻灯片 3：引言 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), GOLD)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(6), Inches(0.7),
             "引  言", font_size=32, font_color=TEXT_WHITE, bold=True)
add_rect(slide, Inches(0.8), Inches(1.1), Inches(2), Inches(0.03), GOLD)

# 内容区域 - 分成两块
# 左块：时代背景
card1 = add_rounded_rect(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(5.2), BG_ACCENT)
add_accent_bar(slide, Inches(0.8), Inches(1.5), Inches(5.2), GOLD)
add_text_box(slide, Inches(1.2), Inches(1.7), Inches(5), Inches(0.5),
             "时代背景", font_size=22, font_color=GOLD, bold=True)

bg_items = [
    "帝制终结（1911）：传统政治秩序的崩塌",
    "西学涌入：科学、民主、进化论冲击传统世界观",
    "新文化运动（1915起）：「打倒孔家店」与思想解放",
    "救亡图存：民族危机下的文化焦虑",
    "「中西文化之争」：中国文化向何处去？",
]
add_bullet_list(slide, Inches(1.2), Inches(2.3), Inches(5), Inches(4),
                bg_items, font_size=16, font_color=TEXT_LIGHT, spacing_after=10)

# 右块：核心问题
card2 = add_rounded_rect(slide, Inches(6.8), Inches(1.5), Inches(5.8), Inches(5.2), BG_ACCENT)
add_accent_bar(slide, Inches(6.8), Inches(1.5), Inches(5.2), ACCENT_TEAL)
add_text_box(slide, Inches(7.2), Inches(1.7), Inches(5), Inches(0.5),
             "核心问题", font_size=22, font_color=ACCENT_TEAL, bold=True)

q_items = [
    "宗教在现代中国还有没有位置？",
    "中国传统文化中的儒释道如何应对西方冲击？",
    "文人如何在「信」与「不信」之间寻找精神安顿？",
    "「科学代宗教」是否可能？",
    "宗教的哲学内核与仪式迷信如何区分？",
]
add_bullet_list(slide, Inches(7.2), Inches(2.3), Inches(5), Inches(4),
                q_items, font_size=16, font_color=TEXT_LIGHT, bullet_color=ACCENT_TEAL, spacing_after=10)

# ==================== 通用内容幻灯片模板 ====================
def make_content_slide(title, subtitle, sections, accent_color=GOLD):
    """创建内容页，sections = [(title, [bullets], color), ...]"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), GOLD)

    add_text_box(slide, Inches(0.8), Inches(0.3), Inches(8), Inches(0.7),
                 title, font_size=30, font_color=TEXT_WHITE, bold=True)
    if subtitle:
        add_text_box(slide, Inches(0.8), Inches(0.95), Inches(8), Inches(0.5),
                     subtitle, font_size=16, font_color=accent_color)
    add_rect(slide, Inches(0.8), Inches(1.35), Inches(2), Inches(0.03), accent_color)

    n = len(sections)
    col_width = (Inches(12)) / n
    start_x = Inches(0.8)

    for idx, (sec_title, bullets, sec_color) in enumerate(sections):
        x = start_x + col_width * idx
        card_h = Inches(5.5)

        # 卡片背景
        add_rounded_rect(slide, x, Inches(1.6), col_width - Inches(0.15), card_h, BG_ACCENT)
        add_accent_bar(slide, x, Inches(1.6), card_h, sec_color)

        # 卡片标题
        add_text_box(slide, x + Inches(0.15), Inches(1.7),
                     col_width - Inches(0.5), Inches(0.5),
                     sec_title, font_size=18, font_color=sec_color, bold=True)

        # 要点列表
        add_bullet_list(slide, x + Inches(0.15), Inches(2.3),
                        col_width - Inches(0.5), Inches(4.5),
                        bullets, font_size=14, font_color=TEXT_LIGHT,
                        bullet_color=sec_color, spacing_after=8)

    return slide

# ==================== 幻灯片 4：一、儒家为体 ====================
make_content_slide(
    "壹  儒家为体：以儒学为信仰",
    "主要特点：视儒学为伦理—文化体系，强调道德实践而非鬼神祭祀",
    [
        ("主要特点", [
            "视儒学为伦理—文化体系，而非神学宗教",
            "强调道德实践而非鬼神祭祀",
            "以「内在超越」替代「外在信仰」",
        ], GOLD),
        ("梁漱溟", [
            "《东西文化及其哲学》（1921）",
            "「儒学是人生哲学，不是宗教」",
            "儒学的「向上心」可替代宗教功能",
            "中国文化走「伦理本位」路线",
        ], ACCENT_RED),
        ("熊十力 & 马一浮", [
            "熊十力：创「新唯识论」，融合儒释",
            "「本心即天」——以道德本心为宇宙本体",
            "马一浮：「六艺统摄一切学术」",
            "融通儒释道三家",
        ], ACCENT_PURPLE),
    ],
    accent_color=GOLD
)

# ==================== 幻灯片 5：二、佛教徒 ====================
make_content_slide(
    "贰  佛教徒：皈依或深度修习",
    "两条路径：信仰实践路径 + 学术研究路径",
    [
        ("主要特点", [
            "因乱世苦痛、生命虚无感而转向佛法",
            "注重解脱、慈悲、因果",
            "倾向于禅宗或净土宗",
            "部分深入唯识学研究",
        ], RGBColor(0x8E, 0x44, 0xAD)),
        ("弘一法师（李叔同）", [
            "从风流才子到律宗高僧",
            "民国最震撼人心的宗教转向",
            "持戒极严，以佛救心",
            "「出家乃大丈夫之事」",
        ], ACCENT_RED),
        ("苏曼殊 & 丰子恺 & 欧阳竟无", [
            "苏曼殊：「情僧」，半僧半俗",
            "丰子恺：《护生画集》，「人间佛教」",
            "欧阳竟无：支那内学院，「以佛学为学问」",
            "两条路径：修行感化 + 学问影响",
        ], ACCENT_TEAL),
    ],
    accent_color=RGBColor(0x8E, 0x44, 0xAD)
)

# ==================== 幻灯片 6：三、基督徒 ====================
make_content_slide(
    "叁  基督徒：信仰或文化认同",
    "文化认同先于信仰认同",
    [
        ("主要特点", [
            "受教会教育影响",
            "视基督教为现代文明象征",
            "重视「爱的哲学」",
            "文化认同先于信仰认同",
        ], RGBColor(0x2C, 0x3E, 0x6B)),
        ("林语堂", [
            "教会学校出身（圣约翰大学）",
            "晚年重归基督教",
            "「回到基督是因为耶稣的人格」",
            "《生活的艺术》融合老庄思想",
        ], ACCENT_ORANGE),
        ("冰心 & 老舍 & 吴雷川", [
            "冰心：基督教博爱思想 →《繁星》《春水》",
            "老舍：未受洗，但对基督教伦理有同情",
            "吴雷川：燕京大学校长，推动基督教本土化",
            "与利玛窦策略相似：以学术/科学为媒介",
        ], ACCENT_GREEN),
    ],
    accent_color=RGBColor(0x2C, 0x3E, 0x6B)
)

# ==================== 幻灯片 7：四、道家/道教 ====================
make_content_slide(
    "肆  道家/道教：隐逸或养生",
    "精神美学 > 宗教信仰",
    [
        ("主要特点", [
            "作为精神姿态，非严格教派归属",
            "追求自然、自由、无为",
            "部分关注养生修炼",
            "追求审美化的「道境」",
        ], ACCENT_TEAL),
        ("周作人", [
            "崇尚「苦茶主义」",
            "冲淡闲适的生活美学",
            "「我只想种一点自己喜欢的东西」",
            "深得老庄「无为」之味",
        ], ACCENT_GREEN),
        ("沈从文 & 朱光潜", [
            "沈从文：信「人性神性」——不在庙堂，在水边",
            "湘西巫傩、泛灵信仰背景",
            "朱光潜：将「逍遥游」转为审美人生态度",
            "「人生的艺术化」= 道家无为的现代诠释",
        ], ACCENT_ORANGE),
    ],
    accent_color=ACCENT_TEAL
)

# ==================== 幻灯片 8：五、无神论者 ====================
make_content_slide(
    "伍  无神论者 / 反宗教者",
    "「科学代宗教」是五四主流叙事",
    [
        ("主要特点", [
            "视宗教为迷信，阻碍科学与启蒙",
            "强调证据的重要性",
            "提倡以科学代替宗教",
            "内部态度有激进与温和之分",
        ], ACCENT_ORANGE),
        ("陈独秀 & 胡适", [
            "陈独秀：「拿证据来」，最激进的反宗教立场",
            "胡适：「宗教是人类幼稚期的产物」",
            "但胡适态度相对温和",
            "承认宗教的道德教化功能",
        ], ACCENT_RED),
        ("鲁迅：需要特别辨析", [
            "猛烈批判民间迷信（《祝福》）",
            "但深入研究佛教——辑校《百喻经》",
            "反的是「迷信」，非宗教哲学内核",
            "更准确：「反迷信者」非「反宗教者」",
        ], ACCENT_PURPLE),
    ],
    accent_color=ACCENT_ORANGE
)

# ==================== 幻灯片 9：六、神秘主义 ====================
make_content_slide(
    "陆  神秘主义 / 泛神论者",
    "以诗性/审美体验替代宗教仪式",
    [
        ("主要特点", [
            "不属于任何教派",
            "相信宇宙有灵性秩序",
            "注重「美与爱即神性」",
            "以诗性体验替代宗教仪式",
        ], RGBColor(0xE7, 0x4C, 0x3C)),
        ("徐志摩", [
            "诗作体现宇宙呼吸与星辉启示",
            "《再别康桥》《偶然》",
            "泛神论式的神秘体验",
            "「自然本身即是神圣的显现」",
        ], GOLD),
        ("废名 & 宗白华", [
            "废名：小说融禅意、童趣、乡土",
            "「不着相」的空灵，接近禅宗",
            "宗白华：美的极致 = 与宇宙生命合一",
            "「天人合一」的宗教性体验",
        ], ACCENT_TEAL),
    ],
    accent_color=RGBColor(0xE7, 0x4C, 0x3C)
)

# ==================== 幻灯片 10：七、学术化研究（新增） ====================
make_content_slide(
    "柒  学术化研究：宗教作为学问",
    "★ 新增章节 ★ 既不信也不反，而是「研究」——民国全新姿态",
    [
        ("主要特点", [
            "宗教从信仰对象转为学术研究对象",
            "运用现代史学、文献学、哲学方法",
            "既不信也不反，保持学术中立",
            "现代宗教学在中国诞生的标志",
        ], ACCENT_GREEN),
        ("汤用彤", [
            "《汉魏两晋南北朝佛教史》（1938）",
            "中国现代佛教史研究的奠基之作",
            "考证文献、辨析源流",
            "从「护教—弘法」中解放出来",
        ], GOLD),
        ("陈寅恪 & 傅斯年", [
            "陈寅恪：精通梵文、巴利文",
            "揭示佛教中国化过程",
            "傅斯年：「史学即史料学」",
            "将宗教文献纳入史料范畴",
        ], ACCENT_PURPLE),
    ],
    accent_color=ACCENT_GREEN
)

# ==================== 幻灯片 11：总结——六个共性 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), GOLD)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(8), Inches(0.7),
             "总结：民国文人宗教观的六个共性", font_size=30, font_color=TEXT_WHITE, bold=True)
add_rect(slide, Inches(0.8), Inches(1.0), Inches(2), Inches(0.03), GOLD)

common_items = [
    ("佛学是公共话语", "无论信不信佛，佛学概念（缘起、空、唯识、顿悟）是共享的思想资源"),
    ("反迷信 ≠ 反宗教", "批判焦点是民间信仰的仪式层面，而非宗教的哲学内核"),
    ("中西比较是基本框架", "几乎所有讨论都放在「中西文化比较」语境中展开"),
    ("实用主义倾向", "关心的不是「宗教是否为真」，而是「宗教对中国有什么用」"),
    ("跨界融合是常态", "儒释道、中西思想融合普遍，纯粹的「单一信仰」反而少见"),
    ("审美化倾向", "宗教体验常被转化为审美体验（禅意、道境、天人合一）——中国独特之处"),
]

for i, (title, desc) in enumerate(common_items):
    y = Inches(1.3) + Inches(i * 0.92)
    colors = [GOLD, ACCENT_RED, ACCENT_TEAL, ACCENT_ORANGE, ACCENT_PURPLE, ACCENT_GREEN]
    color = colors[i]

    # 数字圆圈
    add_number_circle(slide, Inches(1.0), y + Inches(0.05), Inches(0.5), str(i+1), color)

    # 标题
    add_text_box(slide, Inches(1.7), y, Inches(4), Inches(0.4),
                 title, font_size=19, font_color=color, bold=True)
    # 描述
    add_text_box(slide, Inches(1.7), y + Inches(0.4), Inches(10), Inches(0.4),
                 desc, font_size=15, font_color=TEXT_LIGHT)

# ==================== 幻灯片 12：总结——当代启示 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), GOLD)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(8), Inches(0.7),
             "对当代的启示", font_size=30, font_color=TEXT_WHITE, bold=True)
add_rect(slide, Inches(0.8), Inches(1.0), Inches(2), Inches(0.03), GOLD)

# 三个启示卡片
insights = [
    ("启示一", "宗教不仅是信仰问题，更是文化问题",
     "民国文人的态度提醒我们：理解宗教需要超越「信/不信」的二元对立，进入其文化内核。宗教不是非黑即白的选择，而是一个文化系统的深层结构。",
     GOLD),
    ("启示二", "宗教的现代化 ≠ 去宗教化",
     "而是在现代学术框架中重新激活传统宗教思想的解释力——这正是道教理论现代化等当代课题的核心挑战。不是回到古代，而是让传统思想在当代语境中重新发声。",
     ACCENT_TEAL),
    ("启示三", "跨文化对话是宗教研究的必由之路",
     "从利玛窦翻译「太极」「理」到民国文人比较中西宗教，跨文化翻译始终是理解宗教的关键环节。今天的世界更需要这种对话精神。",
     ACCENT_GREEN),
]

for i, (num, title, desc, color) in enumerate(insights):
    x = Inches(0.8) + Inches(i * 4.0)
    w = Inches(3.7)

    card = add_rounded_rect(slide, x, Inches(1.5), w, Inches(5.0), BG_ACCENT)
    add_accent_bar(slide, x, Inches(1.5), Inches(5.0), color)

    add_text_box(slide, x + Inches(0.2), Inches(1.7), w - Inches(0.4), Inches(0.4),
                 num, font_size=16, font_color=color, bold=True)
    add_text_box(slide, x + Inches(0.2), Inches(2.2), w - Inches(0.4), Inches(0.8),
                 title, font_size=20, font_color=TEXT_WHITE, bold=True)
    add_rect(slide, x + Inches(0.2), Inches(3.1), Inches(2), Inches(0.02), color)
    add_text_box(slide, x + Inches(0.2), Inches(3.3), w - Inches(0.4), Inches(2.8),
                 desc, font_size=15, font_color=TEXT_LIGHT)

# ==================== 幻灯片 13：结尾 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), GOLD)
add_rect(slide, Inches(0), Inches(7.44), Inches(13.333), Inches(0.06), GOLD)
add_rect(slide, Inches(0.8), Inches(2.0), Inches(0.06), Inches(3.5), GOLD)

add_text_box(slide, Inches(1.3), Inches(2.2), Inches(10), Inches(1.2),
             "感谢聆听", font_size=48, font_color=TEXT_WHITE, bold=True)

add_rect(slide, Inches(1.3), Inches(3.6), Inches(4), Inches(0.03), GOLD)

add_text_box(slide, Inches(1.3), Inches(4.0), Inches(10), Inches(0.8),
             "宗教不仅是信仰问题，更是文化问题", font_size=22, font_color=GOLD)

add_text_box(slide, Inches(1.3), Inches(5.2), Inches(6), Inches(0.5),
             "戴建华  ｜  2026年6月", font_size=16, font_color=TEXT_LIGHT)

# 保存
output_path = "/home/admin/.openclaw/workspace/民国文人的宗教观.pptx"
prs.save(output_path)
print(f"✅ PPT saved to {output_path}")
print(f"Total slides: {len(prs.slides)}")
