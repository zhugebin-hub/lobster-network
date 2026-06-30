#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
根据图片内容创建可编辑的 PPT 文件 - 六三五精进卓越课程体系
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def create_ppt():
    # 创建演示文稿
    prs = Presentation()
    
    # 使用 16:9 的幻灯片
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # ========== 创建幻灯片 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白版式
    
    # 设置背景色为白色
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    # 左上角标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(3), Inches(0.6))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = "课程体系优化"
    title_para.font.size = Pt(20)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(44, 62, 80)
    
    # 标题左侧红色三角
    triangle = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_TRIANGLE,
        Inches(0.3), Inches(0.55), Inches(0.25), Inches(0.3)
    )
    triangle.fill.solid()
    triangle.fill.fore_color.rgb = RGBColor(231, 76, 60)
    triangle.line.fill.background()
    
    # 顶部红色虚线
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5), Inches(0.95), Inches(12.5), Inches(0.05)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(231, 76, 60)
    line.line.fill.background()
    
    # 主标题横幅
    banner = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(4.5), Inches(0.4), Inches(5), Inches(0.7)
    )
    banner.fill.solid()
    banner.fill.fore_color.rgb = RGBColor(41, 128, 185)
    banner.line.fill.background()
    
    # 横幅文字
    banner_text = slide.shapes.add_textbox(Inches(4.7), Inches(0.55), Inches(4.6), Inches(0.4))
    banner_frame = banner_text.text_frame
    banner_para = banner_frame.paragraphs[0]
    banner_para.text = '"六三五"精进卓越课程体系'
    banner_para.font.size = Pt(18)
    banner_para.font.bold = True
    banner_para.font.color.rgb = RGBColor(255, 255, 255)
    banner_para.alignment = PP_ALIGN.CENTER
    
    # ========== 左侧金字塔图 ==========
    # 基础课程（底层）
    base = slide.shapes.add_shape(
        MSO_SHAPE.TRAPEZOID,
        Inches(0.5), Inches(4.5), Inches(4.5), Inches(1.2)
    )
    base.fill.solid()
    base.fill.fore_color.rgb = RGBColor(52, 152, 219)
    base.line.color.rgb = RGBColor(41, 128, 185)
    base.line.width = Pt(1)
    
    base_text = slide.shapes.add_textbox(Inches(0.5), Inches(4.8), Inches(4.5), Inches(0.6))
    base_text_frame = base_text.text_frame
    base_text_para = base_text_frame.paragraphs[0]
    base_text_para.text = "基础课程"
    base_text_para.font.size = Pt(16)
    base_text_para.font.bold = True
    base_text_para.font.color.rgb = RGBColor(255, 255, 255)
    base_text_para.alignment = PP_ALIGN.CENTER
    
    # 拓展课程（中层）
    middle = slide.shapes.add_shape(
        MSO_SHAPE.TRAPEZOID,
        Inches(1.2), Inches(3.5), Inches(3.1), Inches(1.0)
    )
    middle.fill.solid()
    middle.fill.fore_color.rgb = RGBColor(46, 204, 113)
    middle.line.color.rgb = RGBColor(39, 174, 96)
    middle.line.width = Pt(1)
    
    middle_text = slide.shapes.add_textbox(Inches(1.2), Inches(3.75), Inches(3.1), Inches(0.5))
    middle_text_frame = middle_text.text_frame
    middle_text_para = middle_text_frame.paragraphs[0]
    middle_text_para.text = "拓展课程"
    middle_text_para.font.size = Pt(14)
    middle_text_para.font.bold = True
    middle_text_para.font.color.rgb = RGBColor(255, 255, 255)
    middle_text_para.alignment = PP_ALIGN.CENTER
    
    # 卓越课程（顶层）
    top = slide.shapes.add_shape(
        MSO_SHAPE.ISOSCELES_TRIANGLE,
        Inches(2.0), Inches(2.0), Inches(1.5), Inches(1.5)
    )
    top.fill.solid()
    top.fill.fore_color.rgb = RGBColor(231, 76, 60)
    top.line.color.rgb = RGBColor(192, 57, 43)
    top.line.width = Pt(1)
    
    top_text = slide.shapes.add_textbox(Inches(2.0), Inches(2.5), Inches(1.5), Inches(0.5))
    top_text_frame = top_text.text_frame
    top_text_para = top_text_frame.paragraphs[0]
    top_text_para.text = "卓越课程"
    top_text_para.font.size = Pt(12)
    top_text_para.font.bold = True
    top_text_para.font.color.rgb = RGBColor(255, 255, 255)
    top_text_para.alignment = PP_ALIGN.CENTER
    
    # ========== 右侧三个板块 ==========
    content_boxes = [
        {
            "icon": "🎯",
            "title": "六素养：全面发展的核心",
            "content": "涵盖品格、人文、科学、健康、艺术、劳动六大关键素养，奠定学生成长基石。",
            "x": 5.5,
            "y": 1.5,
            "color": RGBColor(52, 152, 219)
        },
        {
            "icon": "📊",
            "title": "三阶梯：进阶式培养路径",
            "content": "基础课程（共同基础）→ 拓展课程（早期引领）→ 卓越课程（优势潜能）。",
            "x": 5.5,
            "y": 3.2,
            "color": RGBColor(46, 204, 113)
        },
        {
            "icon": "️",
            "title": "五维度：课程实施的保障",
            "content": "国家课程规范化、德育课程主体化、校本课程特色化、社团课程自主化、实践课程多元化。",
            "x": 5.5,
            "y": 4.9,
            "color": RGBColor(142, 68, 173)
        }
    ]
    
    for box in content_boxes:
        # 标题
        title_box = slide.shapes.add_textbox(
            Inches(box["x"]), Inches(box["y"]), Inches(7.5), Inches(0.5)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = box["title"]
        title_para.font.size = Pt(14)
        title_para.font.bold = True
        title_para.font.color.rgb = box["color"]
        
        # 内容
        content_box = slide.shapes.add_textbox(
            Inches(box["x"]), Inches(box["y"] + 0.4), Inches(7.5), Inches(0.8)
        )
        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        content_para = content_frame.paragraphs[0]
        content_para.text = box["content"]
        content_para.font.size = Pt(11)
        content_para.font.color.rgb = RGBColor(80, 80, 80)
    
    # ========== 底部课程分类框 ==========
    # 第一行课程
    course_row1 = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(5.2), Inches(2.8), Inches(0.8), Inches(0.35)
    )
    course_row1.fill.solid()
    course_row1.fill.fore_color.rgb = RGBColor(231, 76, 60)
    course_row1.line.fill.background()
    
    text1 = slide.shapes.add_textbox(Inches(5.25), Inches(2.85), Inches(0.7), Inches(0.25))
    t1 = text1.text_frame.paragraphs[0]
    t1.text = "竞赛培优"
    t1.font.size = Pt(8)
    t1.font.color.rgb = RGBColor(255, 255, 255)
    t1.alignment = PP_ALIGN.CENTER
    
    # 第二行课程
    course_row2 = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(5.2), Inches(3.3), Inches(0.8), Inches(0.35)
    )
    course_row2.fill.solid()
    course_row2.fill.fore_color.rgb = RGBColor(46, 204, 113)
    course_row2.line.fill.background()
    
    text2 = slide.shapes.add_textbox(Inches(5.25), Inches(3.35), Inches(0.7), Inches(0.25))
    t2 = text2.text_frame.paragraphs[0]
    t2.text = "选择性必修"
    t2.font.size = Pt(8)
    t2.font.color.rgb = RGBColor(255, 255, 255)
    t2.alignment = PP_ALIGN.CENTER
    
    # 底部圆形
    circle1 = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(6.5), Inches(5.8), Inches(1.0), Inches(1.0)
    )
    circle1.fill.solid()
    circle1.fill.fore_color.rgb = RGBColor(52, 152, 219)
    circle1.line.fill.background()
    
    c1_text = slide.shapes.add_textbox(Inches(6.6), Inches(6.0), Inches(0.8), Inches(0.6))
    c1 = c1_text.text_frame.paragraphs[0]
    c1.text = "综合实践\n活动"
    c1.font.size = Pt(9)
    c1.font.color.rgb = RGBColor(255, 255, 255)
    c1.alignment = PP_ALIGN.CENTER
    
    circle2 = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(7.8), Inches(5.8), Inches(1.0), Inches(1.0)
    )
    circle2.fill.solid()
    circle2.fill.fore_color.rgb = RGBColor(46, 204, 113)
    circle2.line.fill.background()
    
    c2_text = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.8), Inches(5.8), Inches(1.0), Inches(1.0))
    c2_text.fill.solid()
    c2_text.fill.fore_color.rgb = RGBColor(46, 204, 113)
    c2_text.line.fill.background()
    
    c2_label = slide.shapes.add_textbox(Inches(7.9), Inches(6.1), Inches(0.8), Inches(0.4))
    c2 = c2_label.text_frame.paragraphs[0]
    c2.text = "劳动教育"
    c2.font.size = Pt(9)
    c2.font.color.rgb = RGBColor(255, 255, 255)
    c2.alignment = PP_ALIGN.CENTER
    
    # 保存文件
    output_path = "/home/admin/.openclaw/workspace/六三五精进卓越课程体系.pptx"
    prs.save(output_path)
    print(f"PPT 文件已创建：{output_path}")
    return output_path

if __name__ == "__main__":
    create_ppt()
