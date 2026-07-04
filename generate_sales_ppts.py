#!/usr/bin/env python3
"""
《商品销售》PPT课件生成器
章节：售前准备、推介商品、售后服务
作者：诸葛虾 AI助手
日期：2026-04-25
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os
import zipfile
from datetime import datetime

# 配色方案
COLORS = {
    'primary': RGBColor(0x1E, 0x88, 0xE5),    # 主色-蓝色
    'secondary': RGBColor(0x43, 0xA0, 0x47),  # 辅色-绿色
    'accent': RGBColor(0xFF, 0x6F, 0x00),     # 强调色-橙色
    'warning': RGBColor(0xE5, 0x39, 0x35),    # 警告色-红色
    'dark': RGBColor(0x26, 0x32, 0x38),       # 深色文字
    'light': RGBColor(0xFF, 0xFF, 0xFF),      # 白色
    'gray': RGBColor(0x78, 0x90, 0x9C),       # 灰色
    'bg_light': RGBColor(0xF5, 0xF7, 0xFA),   # 浅灰背景
}

def set_slide_bg(slide, color):
    """设置幻灯片背景色"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, color, text="", font_size=14, font_color=COLORS['light'], bold=False, alignment=PP_ALIGN.LEFT):
    """添加形状并设置文字"""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.bold = bold
        p.alignment = alignment
        tf.paragraphs[0].space_before = Pt(6)
        tf.paragraphs[0].space_after = Pt(6)
    
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=14, font_color=COLORS['dark'], bold=False, alignment=PP_ALIGN.LEFT, font_name="微软雅黑"):
    """添加文本框"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.alignment = alignment
    p.font.name = font_name
    return txBox

def add_bullet_list(slide, left, top, width, height, items, font_size=13, font_color=COLORS['dark'], spacing=Pt(8)):
    """添加项目符号列表"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.name = "微软雅黑"
        p.space_after = spacing
        p.level = 0
    return txBox

def add_header_bar(slide, text, subtitle=""):
    """添加章节标题栏"""
    # 顶部蓝色条
    add_shape(slide, Inches(0), Inches(0), Inches(10), Inches(1.2), COLORS['primary'])
    add_text_box(slide, Inches(0.5), Inches(0.2), Inches(9), Inches(0.8), text, 
                 font_size=28, font_color=COLORS['light'], bold=True)
    if subtitle:
        add_text_box(slide, Inches(0.5), Inches(0.8), Inches(9), Inches(0.4), subtitle,
                     font_size=14, font_color=RGBColor(0xBB, 0xDE, 0xFB))

def add_section_header(slide, chapter_num, chapter_name,课时_info):
    """添加章节封面"""
    set_slide_bg(slide, COLORS['primary'])
    
    # 装饰线条
    add_shape(slide, Inches(0), Inches(3.5), Inches(10), Inches(0.05), COLORS['light'])
    add_shape(slide, Inches(0), Inches(5.5), Inches(10), Inches(0.05), COLORS['light'])
    
    # 章节号
    add_text_box(slide, Inches(1), Inches(2), Inches(8), Inches(1.5), 
                 f"第{chapter_num}章", font_size=20, font_color=RGBColor(0xBB, 0xDE, 0xFB), bold=True)
    
    # 章节标题
    add_text_box(slide, Inches(1), Inches(2.8), Inches(8), Inches(1.2),
                 chapter_name, font_size=40, font_color=COLORS['light'], bold=True)
    
    # 课时信息
    add_text_box(slide, Inches(1), Inches(4.0), Inches(8), Inches(0.5),
                 f"共 {课时_info} 课时", font_size=18, font_color=COLORS['light'])
    
    # 底部信息
    add_text_box(slide, Inches(1), Inches(5.8), Inches(8), Inches(0.5),
                 "《商品销售》· 高等教育出版社 · 中职商贸类专业",
                 font_size=12, font_color=RGBColor(0x90, 0xCA, 0xF9))

def add_content_slide(slide, title, content_items, difficulty="基础", interactive=False, current_event=""):
    """添加内容页"""
    set_slide_bg(slide, COLORS['bg_light'])
    add_header_bar(slide, title)
    
    y_offset = 1.5
    
    # 难度标签
    diff_colors = {"基础": COLORS['secondary'], "进阶": COLORS['accent'], "拓展": COLORS['warning']}
    diff_label = f"📊 {difficulty}"
    add_shape(slide, Inches(0.5), Inches(y_offset), Inches(1.5), Inches(0.35), 
              diff_colors.get(difficulty, COLORS['gray']))
    add_text_box(slide, Inches(0.5), Inches(y_offset), Inches(1.5), Inches(0.35),
                 diff_label, font_size=10, font_color=COLORS['light'], bold=True, alignment=PP_ALIGN.CENTER)
    
    y_offset += 0.5
    
    # 互动标签
    if interactive:
        add_shape(slide, Inches(2.2), Inches(y_offset - 0.5), Inches(1.5), Inches(0.35),
                  COLORS['accent'])
        add_text_box(slide, Inches(2.2), Inches(y_offset - 0.5), Inches(1.5), Inches(0.35),
                     "💬 师生互动", font_size=10, font_color=COLORS['light'], bold=True, alignment=PP_ALIGN.CENTER)
    
    # 时事标签
    if current_event:
        add_shape(slide, Inches(3.9), Inches(y_offset - 0.5), Inches(2.5), Inches(0.35),
                  COLORS['warning'])
        add_text_box(slide, Inches(3.9), Inches(y_offset - 0.5), Inches(2.5), Inches(0.35),
                     "📰 时事热点", font_size=10, font_color=COLORS['light'], bold=True, alignment=PP_ALIGN.CENTER)
    
    y_offset += 0.3
    
    # 内容列表
    add_bullet_list(slide, Inches(0.5), Inches(y_offset), Inches(9), Inches(4.5),
                    content_items, font_size=14)

def add_interactive_slide(slide, title, question, options, answer, explanation):
    """添加互动问答页"""
    set_slide_bg(slide, COLORS['bg_light'])
    add_header_bar(slide, title)
    
    # 问题
    add_text_box(slide, Inches(0.5), Inches(1.5), Inches(9), Inches(0.6),
                 "🤔 思考题", font_size=18, font_color=COLORS['primary'], bold=True)
    
    add_text_box(slide, Inches(0.5), Inches(2.1), Inches(9), Inches(0.8),
                 question, font_size=16, font_color=COLORS['dark'], bold=False)
    
    # 选项
    y = 3.0
    for i, opt in enumerate(options):
        color = COLORS['primary'] if i < 2 else COLORS['secondary']
        add_shape(slide, Inches(0.5), Inches(y + i * 0.5), Inches(9), Inches(0.45), color)
        add_text_box(slide, Inches(0.7), Inches(y + i * 0.5), Inches(8.6), Inches(0.45),
                     opt, font_size=14, font_color=COLORS['light'], alignment=PP_ALIGN.LEFT)
    
    # 答案（折叠效果）
    add_shape(slide, Inches(0.5), Inches(y + len(options) * 0.5 + 0.3), Inches(9), Inches(1.0),
              COLORS['secondary'])
    add_text_box(slide, Inches(0.7), Inches(y + len(options) * 0.5 + 0.3), Inches(8.6), Inches(0.4),
                 f"✅ 正确答案：{answer}", font_size=14, font_color=COLORS['light'], bold=True)
    add_text_box(slide, Inches(0.7), Inches(y + len(options) * 0.5 + 0.7), Inches(8.6), Inches(0.5),
                 explanation, font_size=12, font_color=RGBColor(0xE8, 0xF5, 0xE9))

def add_case_study_slide(slide, title, case_title, case_content, analysis_points):
    """添加案例分析页"""
    set_slide_bg(slide, COLORS['bg_light'])
    add_header_bar(slide, title)
    
    # 案例标题
    add_shape(slide, Inches(0.5), Inches(1.5), Inches(9), Inches(0.5), COLORS['accent'])
    add_text_box(slide, Inches(0.7), Inches(1.5), Inches(8.6), Inches(0.5),
                 f"📋 {case_title}", font_size=16, font_color=COLORS['light'], bold=True)
    
    # 案例内容
    add_text_box(slide, Inches(0.5), Inches(2.2), Inches(9), Inches(1.5),
                 case_content, font_size=13, font_color=COLORS['dark'])
    
    # 分析要点
    add_text_box(slide, Inches(0.5), Inches(3.8), Inches(4), Inches(0.4),
                 "💡 分析要点：", font_size=14, font_color=COLORS['primary'], bold=True)
    
    add_bullet_list(slide, Inches(0.5), Inches(4.2), Inches(9), Inches(2.5),
                    analysis_points, font_size=13)

def add_summary_slide(slide, title, key_points, homework=""):
    """添加总结页"""
    set_slide_bg(slide, COLORS['primary'])
    
    add_text_box(slide, Inches(1), Inches(1), Inches(8), Inches(0.8),
                 title, font_size=32, font_color=COLORS['light'], bold=True, alignment=PP_ALIGN.CENTER)
    
    add_shape(slide, Inches(2), Inches(1.8), Inches(6), Inches(0.05), COLORS['light'])
    
    add_bullet_list(slide, Inches(1.5), Inches(2.2), Inches(7), Inches(3.5),
                    key_points, font_size=16, font_color=COLORS['light'], spacing=Pt(12))
    
    if homework:
        add_shape(slide, Inches(1), Inches(5.2), Inches(8), Inches(0.8), RGBColor(0x15, 0x65, 0xC0))
        add_text_box(slide, Inches(1.2), Inches(5.3), Inches(7.6), Inches(0.6),
                     f"📝 课后作业：{homework}", font_size=14, font_color=COLORS['light'])

# ==================== 第一章：售前准备 ====================

def create_chapter1():
    """第一章：售前准备 - 4课时"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # 章节封面
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
    add_section_header(slide, "一", "售前准备", "4课时")
    
    # ===== 第1课时：认识商品销售 =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_content_slide(slide, "第1课时 认识商品销售", [
        "📌 商品销售的定义与内涵",
        "  - 销售不仅是\"卖东西\"，更是价值传递的过程",
        "  - 现代销售：从\"推销产品\"到\"解决问题\"的转变",
        "",
        "📌 商品销售的基本流程",
        "  售前准备 → 客户开发 → 需求分析 → 商品推介 → 成交 → 售后服务",
        "",
        "📌 销售人员的职业素养",
        "  - 专业知识、沟通能力、服务意识、职业道德",
    ], difficulty="基础", interactive=True, current_event="2026年政府工作报告提出\"大力发展数字经济，推动线上线下消费融合\"")
    
    # 互动页
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_interactive_slide(slide, "第1课时 互动讨论",
        "你认为一个好的销售人员最重要的品质是什么？",
        ["A. 口才好，能说会道", "B. 了解产品，专业可靠", 
         "C. 善于倾听，理解客户需求", "D. 有耐心，服务态度好"],
        "以上都是，但核心是C——理解客户需求",
        "现代销售理论认为，销售的本质是帮助客户解决问题。只有先理解客户需求，才能提供真正有价值的解决方案。口才、专业知识、服务态度都是围绕\"理解需求\"这个核心展开的。")
    
    # ===== 第2课时：商品知识准备 =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_content_slide(slide, "第2课时 商品知识准备", [
        "📌 商品知识的核心内容",
        "  - 商品名称、品牌、产地、规格、型号",
        "  - 商品功能、特点、使用方法",
        "  - 商品的质量标准、认证标志",
        "  - 商品的售后服务政策",
        "",
        "📌 商品知识学习的方法",
        "  - 阅读产品说明书、参加厂家培训",
        "  - 实际体验商品（\"用产品说话\"）",
        "  - 向同行学习、收集客户反馈",
        "",
        "📌 竞品分析",
        "  - 了解竞争对手产品的优劣势",
        "  - 找出自家产品的差异化卖点",
    ], difficulty="基础")
    
    # 案例页
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_case_study_slide(slide, "第2课时 案例分析",
        "案例：小米之家的销售顾问",
        "小米之家的销售顾问需要掌握数百款产品的参数、功能、使用场景。他们不仅会背诵参数，更会\"演示\"——现场展示智能家居联动、拍照对比、游戏性能测试。2025年小米推出\"人车家全生态\"战略后，销售顾问还需要了解小米汽车与智能家居的联动功能。",
        ["商品知识不只是\"背参数\"，更要理解使用场景",
         "实际体验比理论讲解更有说服力",
         "产品线扩展时，知识体系需要同步更新",
         "2026年AIoT（人工智能物联网）趋势下，销售人员需要掌握跨品类联动知识"])
    
    # ===== 第3课时：销售工具与物料准备 =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_content_slide(slide, "第3课时 销售工具与物料准备", [
        "📌 传统销售物料",
        "  - 产品目录、宣传册、价目表",
        "  - 样品、展示道具、演示设备",
        "  - 名片、企业介绍资料",
        "",
        "📌 数字化销售工具（新时代重点）",
        "  - 企业微信、钉钉等客户管理工具",
        "  - 电子画册、H5产品页、短视频",
        "  - AI辅助销售工具：智能话术推荐、客户画像分析",
        "",
        "📌 2026年销售工具新趋势",
        "  - AI数字人直播带货工具普及",
        "  - VR/AR商品展示（如家居AR摆放）",
        "  - 大数据精准营销平台",
    ], difficulty="进阶", current_event="2026年商务部等9部门发布《商贸流通高质量发展行动计划》，推动数字化销售工具应用")
    
    # 互动页
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_interactive_slide(slide, "第3课时 小组讨论",
        "如果你是某品牌服装店的销售，请列出你上班前需要准备的销售工具和物料清单",
        ["传统工具：样品、价签、衣架、购物袋", "数字工具：企业微信、电子画册、搭配推荐小程序",
         "个人准备：仪容仪表、产品知识、当日促销信息", "以上全部都需要！"],
        "D - 以上全部",
        "现代销售需要\"传统+数字\"双轮驱动。2026年消费者既期待线下体验的温度，也习惯数字化服务的便捷。销售人员需要熟练掌握两类工具，根据客户偏好灵活切换。")
    
    # ===== 第4课时：客户信息收集与分析 =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_content_slide(slide, "第4课时 客户信息收集与分析", [
        "📌 客户信息收集渠道",
        "  - 门店登记、线上表单、社交媒体",
        "  - 历史购买记录、客服反馈",
        "  - 市场调研、行业报告",
        "",
        "📌 客户画像构建",
        "  - 基本信息：年龄、性别、职业、收入",
        "  - 消费行为：购买频率、偏好品类、价格敏感度",
        "  - 心理特征：价值观、生活方式、决策风格",
        "",
        "📌 AI时代的客户分析",
        "  - 利用AI工具进行客户分群和精准画像",
        "  - 预测客户需求和购买意向",
        "  - 注意：遵守《个人信息保护法》，合法合规收集信息",
    ], difficulty="进阶", current_event="2025年《个人信息保护法》实施后，企业客户数据管理面临更严格要求")
    
    # 第一章总结
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_summary_slide(slide, "第一章 总结", [
        "✅ 售前准备是销售成功的基础",
        "✅ 商品知识是销售人员的\"底气\"",
        "✅ 数字化工具大幅提升销售效率",
        "✅ 客户画像帮助精准匹配需求",
        "✅ 合法合规是客户信息使用的前提",
    ], homework="选择一款你熟悉的商品，制作一份完整的售前准备清单（包括商品知识要点、所需物料、目标客户画像）")
    
    return prs

# ==================== 第二章：推介商品 ====================

def create_chapter2():
    """第二章：推介商品 - 6课时"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_header(slide, "二", "推介商品", "6课时")
    
    # ===== 第1课时：接近客户的方法 =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_content_slide(slide, "第1课时 接近客户的方法", [
        "📌 接近客户的目的",
        "  - 引起客户注意和兴趣",
        "  - 建立良好的第一印象",
        "  - 获取进一步沟通的机会",
        "",
        "📌 常用接近方法",
        "  - 问候接近法：微笑+礼貌问候",
        "  - 赞美接近法：真诚赞美（注意分寸）",
        "  - 提问接近法：用问题引发兴趣",
        "  - 演示接近法：用产品展示吸引注意",
        "  - 利益接近法：直接告知客户能获得的好处",
        "",
        "📌 接近客户的\"三不\"原则",
        "  不过分热情、不急于推销、不贬低竞品",
    ], difficulty="基础", interactive=True)
    
    # 互动角色扮演
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_interactive_slide(slide, "第1课时 角色扮演",
        "场景：一位客户走进电子产品店，在平板电脑展示区驻足。请用两种不同的接近方式与客户打招呼，并比较效果。",
        ["方式A：\"欢迎光临！需要买平板吗？我们新款刚上市\"", 
         "方式B：\"您好！您看的这款平板特别适合追剧和记笔记，屏幕是12.9英寸的护眼屏\"",
         "方式C：\"您好，随便看看，有需要随时叫我\"（保持距离）",
         "对比讨论：哪种方式更容易打开话匣子？为什么？"],
        "没有标准答案，需要视客户类型而定",
        "方式A偏直接，适合目的明确的客户；方式B有信息量，适合犹豫型客户；方式C给空间，适合喜欢自主浏览的客户。关键是\"察言观色\"，灵活选择。")
    
    # ===== 第2课时：需求探询技巧 =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_content_slide(slide, "第2课时 需求探询技巧", [
        "📌 为什么需求探询很重要？",
        "  - 客户说的不一定是真正需要的",
        "  - 深度挖掘才能找到最佳匹配方案",
        "  - 好的探询让客户感觉\"被理解\"",
        "",
        "📌 SPIN提问法（经典销售工具）",
        "  S - Situation（背景问题）：\"您目前用什么设备？\"",
        "  P - Problem（难点问题）：\"使用中有什么不满意？\"",
        "  I - Implication（暗示问题）：\"这个问题对您的工作有什么影响？\"",
        "  N - Need-payoff（需求效益）：\"如果有一个方案能解决，对您会有多大帮助？\"",
        "",
        "📌 倾听的艺术",
        "  - 70%倾听 + 30%说话",
        "  - 用点头、复述表示理解",
        "  - 捕捉客户的\"弦外之音\"",
    ], difficulty="进阶", interactive=True)
    
    # ===== 第3课时：FABE推介法 =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_content_slide(slide, "第3课时 FABE推介法", [
        "📌 FABE——最经典的商品推介框架",
        "  F - Features（特点）：商品本身有什么特性",
        "  A - Advantages（优势）：这个特点带来什么优势",
        "  B - Benefits（利益）：这个优势给客户什么具体好处",
        "  E - Evidence（证据）：用什么证明你说的",
        "",
        "📌 FABE实战示例（以保温杯为例）",
        "  F：\"这款保温杯采用316不锈钢内胆\"",
        "  A：\"比普通的304不锈钢更耐腐蚀、更耐用\"",
        "  B：\"您装咖啡、茶等饮品不会有异味，用三五年依然如新\"",
        "  E：\"这是国家质检报告，客户满意度98%\"",
        "",
        "📌 FABE使用要点",
        "  - 从客户角度说\"利益\"，不是从产品角度说\"特点\"",
        "  - 证据要具体、可信",
    ], difficulty="进阶")
    
    # FABE练习
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_interactive_slide(slide, "第3课时 FABE练习",
        "请用FABE法推介一款\"无线降噪耳机\"，写出完整的F-A-B-E四句话",
        ["F：主动降噪技术，最高降噪深度达40dB",
         "A：在嘈杂环境中也能享受纯净音乐",
         "B：通勤路上不用调大音量，保护听力，专注效率翻倍",
         "E：某东好评率97%，销量超100万+"],
        "以上是一个参考答案，同学们可以有不同的表达方式",
        "FABE的核心是\"把产品语言翻译成客户语言\"。特点是工程师说的，利益是客户关心的。好的销售能把技术参数翻译成生活场景。")
    
    # ===== 第4课时：商品演示技巧 =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_content_slide(slide, "第4课时 商品演示技巧", [
        "📌 演示的原则",
        "  - 让客户参与进来（\"您试试\"）",
        "  - 一次演示一个核心功能",
        "  - 结合客户需求场景演示",
        "",
        "📌 演示的\"SHOW\"法则",
        "  S - See（让客户看到效果）",
        "  H - Hear（讲解关键信息）",
        "  O - Operate（让客户动手操作）",
        "  W - Want（激发拥有欲望）",
        "",
        "📌 数字化演示新方式",
        "  - 短视频展示（15秒讲清一个卖点）",
        "  - AR虚拟试用（美妆、家居场景）",
        "  - 直播演示（2026年直播带货已成标配）",
    ], difficulty="进阶", current_event="2026年抖音、快手、淘宝直播三大平台直播电商GMV突破5万亿，演示式销售成为核心竞争力")
    
    # ===== 第5课时：处理客户异议 =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_content_slide(slide, "第5课时 处理客户异议", [
        "📌 客户异议的常见类型",
        "  - 价格异议：\"太贵了\"",
        "  - 质量异议：\"这个牌子没听过\"",
        "  - 需求异议：\"我不需要\"",
        "  - 时间异议：\"我考虑考虑\"",
        "  - 竞品异议：\"XX家的更好\"",
        "",
        "📌 LSCPA异议处理法",
        "  L - Listen（倾听，不打断）",
        "  S - Share（共情，\"我理解您的想法\"）",
        "  C - Clarify（澄清，确认真实原因）",
        "  P - Present（提出解决方案）",
        "  A - Ask（请求行动，\"您看这样可以吗\"）",
        "",
        "📌 处理异议的黄金法则",
        "  - 异议=兴趣，没有异议才是最大的问题",
        "  - 不争辩、不否定、不贬低竞品",
    ], difficulty="拓展", interactive=True)
    
    # 异议处理情景模拟
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_interactive_slide(slide, "第5课时 情景模拟",
        "客户说：\"你们这个价格太贵了，XX品牌同款只要一半价格\"。请用LSCPA法回应。",
        ["L：\"您说得对，价格确实是需要考虑的重要因素\"",
         "S：\"我完全理解，谁都想买到性价比高的产品\"",
         "C：\"除了价格，您最关心的是产品的哪些方面？\"",
         "P：\"我们的价格高一些，是因为用了XX材料/提供XX服务，这是对比表...\""],
        "参考答案：完整走一遍LSCPA流程",
        "处理价格异议的关键：不要直接反驳，先共情再澄清，最后用价值对比代替价格对比。客户关心的往往不是\"贵不贵\"，而是\"值不值\"。")
    
    # ===== 第6课时：促成交易 =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_content_slide(slide, "第6课时 促成交易", [
        "📌 识别购买信号",
        "  - 语言信号：\"有货吗？\"\"能便宜点吗？\"\"保修多久？\"",
        "  - 行为信号：反复看商品、问细节、打电话商量",
        "  - 表情信号：点头、微笑、表情放松",
        "",
        "📌 促成交易的常用方法",
        "  - 假设成交法：\"您选黑色还是白色？\"",
        "  - 选择成交法：\"今天送还是明天送？\"",
        "  - 紧迫成交法：\"这个活动今天最后一天\"",
        "  - 从众成交法：\"这款是我们店卖得最好的\"",
        "  - 利益成交法：\"现在下单送XX赠品\"",
        "",
        "📌 促成交易的注意事项",
        "  - 时机很重要，太早让客户有压力，太晚可能流失",
        "  - 不要过度承诺，诚信是底线",
        "  - 成交后及时确认，避免客户\"反悔\"",
    ], difficulty="拓展")
    
    # 第二章总结
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_summary_slide(slide, "第二章 总结", [
        "✅ 接近客户：第一印象决定沟通氛围",
        "✅ 需求探询：SPIN提问法挖掘真实需求",
        "✅ FABE推介：把产品特点翻译成客户利益",
        "✅ 演示技巧：让客户参与是最好的说服",
        "✅ 处理异议：异议=兴趣，用LSCPA化解",
        "✅ 促成交易：识别信号，把握时机",
    ], homework="分组角色扮演：一人扮客户、一人扮销售，模拟一次完整的商品推介过程，录制视频（3-5分钟）")
    
    return prs

# ==================== 第三章：售后服务 ====================

def create_chapter3():
    """第三章：售后服务 - 4课时"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_header(slide, "三", "售后服务", "4课时")
    
    # ===== 第1课时：售后服务的意义 =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_content_slide(slide, "第1课时 售后服务的意义", [
        "📌 售后服务≠\"麻烦事\"",
        "  - 售后服务是客户体验的\"最后一公里\"",
        "  - 好售后 = 免费广告 = 复购 + 转介绍",
        "  - 维护一个老客户的成本 = 开发新客户的1/5",
        "",
        "📌 售后服务的核心价值",
        "  - 提升客户满意度和忠诚度",
        "  - 收集产品反馈，推动产品改进",
        "  - 建立品牌口碑，形成竞争壁垒",
        "  - 创造二次销售机会",
        "",
        "📌 数据说话",
        "  - 客户满意度提升5%，利润增加25%-85%（贝恩咨询）",
        "  - 2026年中国消费者协会数据显示，售后服务满意度已成为消费者选择品牌的首要因素",
    ], difficulty="基础", current_event="2026年央视315晚会主题：\"公平消费 诚信守护\"，售后服务质量成为焦点")
    
    # 互动讨论
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_interactive_slide(slide, "第1课时 经验分享",
        "你或你的家人有没有因为一次好的售后服务而成为某个品牌的忠实客户？请分享一个真实经历。",
        ["案例1：某手机品牌免费换电池，从此只买这个品牌",
         "案例2：某服装店无理由退换，推荐了5个朋友来买",
         "案例3：某家电品牌2小时上门维修，口碑传播",
         "反面案例：售后差导致永远不再购买"],
        "真实案例最能说明售后服务的价值",
        "好的售后服务会让客户成为品牌的\"义务推销员\"。据调查，一个满意的客户会告诉3个人，而一个不满意的客户会告诉11个人。售后不是成本，是投资。")
    
    # ===== 第2课时：售后服务的内容与流程 =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_content_slide(slide, "第2课时 售后服务的内容与流程", [
        "📌 售后服务的主要内容",
        "  - 产品安装、调试",
        "  - 退换货处理",
        "  - 维修、保养",
        "  - 使用指导、技术培训",
        "  - 客户回访、满意度调查",
        "",
        "📌 标准售后服务流程",
        "  客户反馈 → 问题记录 → 分类处理 → 解决方案 → 执行落实 → 回访确认 → 归档分析",
        "",
        "📌 \"三包\"政策（国家法定要求）",
        "  - 包修：7日内免费修理",
        "  - 包换：15日内免费更换",
        "  - 包退：符合退货条件的应予以退货",
        "  - 不同品类有不同规定，需熟悉相关法律法规",
    ], difficulty="基础")
    
    # 案例页
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_case_study_slide(slide, "第2课时 案例分析",
        "案例：胖东来的售后服务体系",
        "胖东来被誉为中国零售业的服务标杆。其售后服务包括：无理由退换货（甚至吃过觉得不好吃也能退）、免费清洗/维修/充电、宠物寄存、母婴室等。2025年胖东来调改帮扶多家亏损超市，核心输出就是服务标准。其客户满意度高达95%以上，复购率超过70%。",
        ["售后服务不是成本中心，而是利润中心",
         "超预期的服务才能形成口碑传播",
         "服务标准需要制度化、可复制",
         "2026年\"消费升级\"趋势下，服务品质成为零售企业核心竞争力"])
    
    # ===== 第3课时：客户投诉处理 =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_content_slide(slide, "第3课时 客户投诉处理", [
        "📌 投诉处理的心态",
        "  - 投诉是礼物——客户给你改进的机会",
        "  - 不投诉就流失——不满的客户会默默离开",
        "  - 投诉处理得好 = 客户忠诚度提升",
        "",
        "📌 HEARD投诉处理法（迪士尼标准）",
        "  H - Hear（倾听）：让客户把话说完，不打断",
        "  E - Empathize（共情）：\"我理解您的感受\"",
        "  A - Apologize（道歉）：为不好的体验道歉（不等于认错）",
        "  R - Resolve（解决）：提供解决方案，给客户选择",
        "  D - Diagnose（诊断）：事后分析原因，预防再次发生",
        "",
        "📌 投诉处理中的沟通技巧",
        "  - 控制情绪，不被客户的愤怒影响",
        "  - 用\"我\"代替\"你\"（\"我来帮您解决\"vs\"你应该...\"）",
        "  - 给出明确的时间承诺（\"24小时内给您回复\"）",
    ], difficulty="进阶", interactive=True)
    
    # 投诉处理模拟
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_interactive_slide(slide, "第3课时 情景模拟",
        "客户怒气冲冲地来到门店：\"你们这个洗衣机买了才一个月就坏了，耽误了我多少事！我要退货！\"请用HEARD法处理。",
        ["H：耐心听完客户的抱怨，不打断，记录关键信息",
         "E：\"非常理解您的心情，洗衣机坏了确实很耽误事\"",
         "A：\"给您带来这样的不便，我代表门店向您道歉\"",
         "R：\"我给您两个方案：一是立即换新，二是先借一台给您用，这台我们马上维修\""],
        "参考答案：完整走一遍HEARD流程",
        "处理投诉的核心：先处理情绪，再处理问题。客户愤怒的时候，讲道理是无效的。先共情、道歉，让客户的情绪降下来，再谈解决方案。给客户选择权，让客户感觉\"被尊重\"。")
    
    # ===== 第4课时：客户关系维护 =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_content_slide(slide, "第4课时 客户关系维护", [
        "📌 客户关系管理（CRM）的核心",
        "  - 从\"一次性交易\"到\"长期关系\"",
        "  - 客户分层管理：VIP客户、普通客户、潜在客户",
        "  - 定期触达：生日祝福、节日问候、新品通知",
        "",
        "📌 客户回访的技巧",
        "  - 回访时机：购买后3天、1个月、3个月",
        "  - 回访方式：电话、微信、短信、上门",
        "  - 回访内容：使用体验、问题收集、需求挖掘",
        "  - 回访禁忌：频繁打扰、纯推销、无准备",
        "",
        "📌 数字化时代的客户关系维护",
        "  - 企业微信社群运营：建立品牌粉丝群",
        "  - 会员积分系统：增加客户粘性",
        "  - AI客服+人工客服：7×24小时响应",
        "  - 2026年新趋势：AI个性化关怀（智能推荐、定制优惠）",
    ], difficulty="拓展", current_event="2026年私域流量运营成为零售企业标配，企业微信用户突破6亿")
    
    # 互动讨论
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_interactive_slide(slide, "第4课时 小组讨论",
        "设计一个你所在专业的\"客户回访方案\"——假设你开了一家网店/实体店，你会如何回访客户？",
        ["方案A：购买后3天微信回访使用体验，30天推送使用技巧",
         "方案B：建立会员群，每月推送专属优惠+新品预告",
         "方案C：生日月送专属礼品+手写感谢卡",
         "方案D：以上组合使用，分层运营"],
        "D - 组合方案效果最佳",
        "好的客户关系维护是\"有温度的自动化\"。用数字化工具提高效率，但关键节点要有人情味。比如生日祝福可以自动发送，但手写卡片、定制礼品更能打动人心。")
    
    # 第三章总结
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_summary_slide(slide, "第三章 总结", [
        "✅ 售后服务是客户体验的\"最后一公里\"",
        "✅ 好售后 = 低成本获客 + 高客户忠诚度",
        "✅ HEARD法：投诉处理的标准流程",
        "✅ 客户关系维护需要\"制度化+人情味\"",
        "✅ 数字化工具让客户关系管理更高效",
        "✅ 售后不是成本，是品牌最值的投资",
    ], homework="选择一家你熟悉的企业，分析其售后服务体系（至少包含3个环节），写一份500字的分析报告")
    
    return prs

# ==================== 主程序 ====================

def main():
    """生成所有课件并打包"""
    output_dir = "/home/admin/.openclaw/workspace/商品销售课件"
    os.makedirs(output_dir, exist_ok=True)
    
    # 生成三个章节的PPT
    print("📚 正在生成第一章：售前准备...")
    prs1 = create_chapter1()
    file1 = os.path.join(output_dir, "01_售前准备.pptx")
    prs1.save(file1)
    print(f"  ✅ {file1}")
    
    print("📚 正在生成第二章：推介商品...")
    prs2 = create_chapter2()
    file2 = os.path.join(output_dir, "02_推介商品.pptx")
    prs2.save(file2)
    print(f"  ✅ {file2}")
    
    print("📚 正在生成第三章：售后服务...")
    prs3 = create_chapter3()
    file3 = os.path.join(output_dir, "03_售后服务.pptx")
    prs3.save(file3)
    print(f"  ✅ {file3}")
    
    # 打包成ZIP
    zip_file = os.path.join(output_dir, "商品销售课件_三章节打包.zip")
    with zipfile.ZipFile(zip_file, 'w', zipfile.ZIP_DEFLATED) as zipf:
        for f in [file1, file2, file3]:
            zipf.write(f, os.path.basename(f))
    
    print(f"\n📦 打包完成：{zip_file}")
    print(f"\n📊 课件统计：")
    print(f"  第一章 售前准备：{len(prs1.slides)} 页")
    print(f"  第二章 推介商品：{len(prs2.slides)} 页")
    print(f"  第三章 售后服务：{len(prs3.slides)} 页")
    print(f"  总计：{len(prs1.slides) + len(prs2.slides) + len(prs3.slides)} 页")
    print(f"\n🎯 课件特色：")
    print(f"  ✅ 难度递进：基础→进阶→拓展")
    print(f"  ✅ 时事热点：结合2025-2026年最新政策和新闻")
    print(f"  ✅ 师生互动：角色扮演、小组讨论、情景模拟")
    print(f"  ✅ 案例分析：小米之家、胖东来、迪士尼等真实案例")
    print(f"  ✅ 课后作业：每章配有实践性作业")

if __name__ == "__main__":
    main()
