#!/usr/bin/env python3
from pptx import Presentation
from pptx.util import Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

prs = Presentation()
prs.slide_width = Cm(33.866)
prs.slide_height = Cm(19.05)
DB=RGBColor(0,51,102); MB=RGBColor(0,102,153); LB=RGBColor(200,230,255)
WH=RGBColor(255,255,255); BK=RGBColor(0,0,0); GY=RGBColor(128,128,128)
LG=RGBColor(240,240,240); RD=RGBColor(204,0,0); TS=15

def sf(r,s=Pt(18),c=BK,b=False):
    r.font.size=s; r.font.color.rgb=c; r.font.bold=b
    rPr = r._r.get_or_add_rPr()
    rPr.set(qn('w:eastAsia'),'微软雅黑')

def T(s,l,t,w,h,tx,fs=14,c=BK,b=False,a=PP_ALIGN.LEFT):
    x=s.shapes.add_textbox(l,t,w,h); f=x.text_frame; f.word_wrap=True; p=f.paragraphs[0]; p.alignment=a
    r=p.add_run(); sf(r,Pt(fs),c,b); r.text=tx; return x

def S(s,l,t,w,h,fc,tx='',fs=12,fc2=BK,b=False,a=PP_ALIGN.CENTER,lc=None):
    x=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,l,t,w,h); x.fill.solid(); x.fill.fore_color.rgb=fc
    if lc: x.line.color.rgb=lc; x.line.width=Pt(1.5)
    else: x.line.fill.background()
    if tx:
        f=x.text_frame; f.word_wrap=True; p=f.paragraphs[0]; p.alignment=a
        r=p.add_run(); sf(r,Pt(fs),fc2,b); r.text=tx
    return x

def R(s,l,t,w,h,fc,tx='',fs=12,fc2=BK,b=False,a=PP_ALIGN.CENTER,lc=None):
    x=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,l,t,w,h); x.fill.solid(); x.fill.fore_color.rgb=fc
    if lc: x.line.color.rgb=lc; x.line.width=Pt(2)
    else: x.line.fill.background()
    if tx:
        f=x.text_frame; f.word_wrap=True; p=f.paragraphs[0]; p.alignment=a
        r=p.add_run(); sf(r,Pt(fs),fc2,b); r.text=tx
    return x

def IP(s,l,t,w,h,lb='📷 图片占位'):
    return S(s,l,t,w,h,LG,lb,14,GY,False,PP_ALIGN.CENTER,RGBColor(180,180,180))

def TB(s,tx):
    S(s,Cm(0),Cm(0),Cm(33.866),Cm(2.2),DB)
    T(s,Cm(1.5),Cm(0.4),Cm(30),Cm(1.5),tx,24,WH,True,PP_ALIGN.LEFT)
    S(s,Cm(0),Cm(2.2),Cm(33.866),Cm(0.06),MB)

def SN(s,n): T(s,Cm(30),Cm(18.2),Cm(3.5),Cm(0.6),f'{n} / {TS}',10,GY,False,PP_ALIGN.RIGHT)
def FT(s): T(s,Cm(1),Cm(18.2),Cm(10),Cm(0.6),'QPPB技术详解 | 芦熠檑 | 2026',9,GY)

# S1 Cover
s=prs.slides.add_slide(prs.slide_layouts[6])
S(s,Cm(0),Cm(0),Cm(33.866),Cm(0.3),DB)
T(s,Cm(3),Cm(5.5),Cm(28),Cm(3),'QPPB技术详解',40,DB,True,PP_ALIGN.CENTER)
T(s,Cm(3),Cm(9),Cm(28),Cm(1.5),'—— 基于 BGP 路由属性的 QoS 策略传播机制 ——',20,MB,False,PP_ALIGN.CENTER)
S(s,Cm(12),Cm(11),Cm(10),Cm(0.05),MB)
T(s,Cm(8),Cm(12.5),Cm(18),Cm(1.5),'汇报人：芦熠檑',16,GY,False,PP_ALIGN.CENTER)
T(s,Cm(8),Cm(14),Cm(18),Cm(1.5),'日期：2026 年 5 月',16,GY,False,PP_ALIGN.CENTER)

# S2 TOC
s=prs.slides.add_slide(prs.slide_layouts[6]); TB(s,'目  录'); SN(s,2); FT(s)
for i,(n,t) in enumerate([('01','QPPB 概述与产生背景'),('02','QPPB 实现原理'),('03','QPPB 典型应用场景')]):
    y=Cm(4+i*3.5); T(s,Cm(4),y,Cm(3),Cm(2),n,36,DB,True)
    T(s,Cm(7.5),y+Cm(0.3),Cm(20),Cm(1.5),t,22,BK,False)
    S(s,Cm(4),y+Cm(2.5),Cm(26),Cm(0.03),LB)

# S3 Part1 Title
s=prs.slides.add_slide(prs.slide_layouts[6])
bg=s.background; bg.fill.solid(); bg.fill.fore_color.rgb=DB
T(s,Cm(4),Cm(7),Cm(26),Cm(3),'第一部分',36,WH,True,PP_ALIGN.CENTER)
T(s,Cm(4),Cm(10.5),Cm(26),Cm(2),'QPPB 概述与产生背景',24,LB,False,PP_ALIGN.CENTER); SN(s,3)

# S4 What is QPPB
s=prs.slides.add_slide(prs.slide_layouts[6]); TB(s,'1.1 什么是 QPPB？'); SN(s,4); FT(s)
items=[('全称','QoS Policy Propagation through BGP\n（通过 BGP 传播 QoS 策略）'),
       ('定义','一种特殊的复杂流分类方法，\n通过 BGP 路由属性对报文进行流分类'),
       ('核心思想','BGP 路由发送者设置 BGP 属性\n预先分类，接收端自动匹配\n关联 QoS 策略'),
       ('优势','网络结构变化时，只需修改\n发送端配置，接收端无需改动')]
for i,(lb,tx) in enumerate(items):
    y=Cm(3+i*3.5); S(s,Cm(2),y,Cm(3),Cm(1.2),DB,lb,14,WH,True)
    T(s,Cm(5.5),y+Cm(0.1),Cm(25),Cm(1.2),tx,14,BK,False)

# S5 Background
s=prs.slides.add_slide(prs.slide_layouts[6]); TB(s,'1.2 产生背景（痛点分析）'); SN(s,5); FT(s)
T(s,Cm(2),Cm(3),Cm(14),Cm(1.2),'传统方式的痛点',16,RD,True)
probs=['AS400 是高优先级网络，需对往返报文重新设置 IP Precedence',
       'Node-A/B 需针对大量 IP 地址配置流分类',
       '网络结构不稳定时，配置修改工作量巨大',
       '大量流分类规则难以维护']
for i,p in enumerate(probs):
    T(s,Cm(2.5),Cm(4.8+i*1.8),Cm(13),Cm(1.5),f'• {p}',12,BK,False)
T(s,Cm(18),Cm(3),Cm(14),Cm(1.2),'QPPB 解决方案',16,DB,True)
sols=['按 AS 信息、团体属性等聚类分类',
      '发送端设置属性，接收端自动匹配',
      '网络变化只需修改发送端配置',
      '大幅简化配置和维护工作量']
for i,p in enumerate(sols):
    T(s,Cm(18.5),Cm(4.8+i*1.8),Cm(13),Cm(1.5),f'✓ {p}',12,BK,False)
S(s,Cm(16.5),Cm(3),Cm(0.03),Cm(12),LB)
IP(s,Cm(2),Cm(13),Cm(30),Cm(5),'📷 图 1-52 跨 AS 组网示意图')

# S6 Positioning
s=prs.slides.add_slide(prs.slide_layouts[6]); TB(s,'1.3 QPPB技术定位'); SN(s,6); FT(s)
td=[['对比维度','简单流分类','QPPB（复杂流分类）'],
    ['分类依据','报文头部固定字段\n（DSCP / 802.1p）','BGP 路由属性\n（AS_PATH / Community）'],
    ['配置复杂度','低（直接基于报文头部）','中（需配置 BGP 属性）'],
    ['适用场景','单域/简单网络','跨 AS / 大型复杂组网'],
    ['灵活性','低（网络变化需重新配置）','高（发送端修改即可）'],
    ['维护成本','高（需逐设备配置）','低（接收端自动适配）']]
tbl=s.shapes.add_table(len(td),3,Cm(2),Cm(3),Cm(30),Cm(12)).table
tbl.columns[0].width=Cm(5); tbl.columns[1].width=Cm(12.5); tbl.columns[2].width=Cm(12.5)
for ri,rd in enumerate(td):
    for ci,ct in enumerate(rd):
        c=tbl.cell(ri,ci); c.text=''; p=c.text_frame.paragraphs[0]; r=p.add_run(); sf(r,12,BK,False); r.text=ct
        p.alignment=PP_ALIGN.CENTER; c.vertical_anchor=MSO_ANCHOR.MIDDLE
        if ri==0: c.fill.solid(); c.fill.fore_color.rgb=DB
        elif ri%2==0: c.fill.solid(); c.fill.fore_color.rgb=LG

# S7 Part2 Title
s=prs.slides.add_slide(prs.slide_layouts[6])
bg=s.background; bg.fill.solid(); bg.fill.fore_color.rgb=DB
T(s,Cm(4),Cm(7),Cm(26),Cm(3),'第二部分',36,WH,True,PP_ALIGN.CENTER)
T(s,Cm(4),Cm(10.5),Cm(26),Cm(2),'QPPB 实现原理',24,LB,False,PP_ALIGN.CENTER); SN(s,7)

# S8 Workflow
s=prs.slides.add_slide(prs.slide_layouts[6]); TB(s,'2.1 QPPB 工作流程（图 1-53）'); SN(s,8); FT(s)
steps=[('步骤 1','BGP 路由发送者\n（Node-C）为路由\n设置特定属性\n（AS_PATH / Community）'),
       ('步骤 2','BGP 路由携带属性\n在 AS 间通告\n（属性作为分类标识）'),
       ('步骤 3','接收者（Node-A）\n匹配路由属性\n设置 Behavior ID\n到 FIB 表'),
       ('步骤 4','数据转发时\n根据目的网络从 FIB\n获取 Behavior ID\n执行对应流动作')]
for i,(lb,tx) in enumerate(steps):
    x=Cm(1.5+i*8)
    R(s,x,Cm(3),Cm(7),Cm(7),WH,lb+'\n\n'+tx,11,BK,False,PP_ALIGN.CENTER,DB)
    if i<3:
        ar=s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,x+Cm(7.2),Cm(5.5),Cm(0.8),Cm(1))
        ar.fill.solid(); ar.fill.fore_color.rgb=MB; ar.line.fill.background()
IP(s,Cm(2),Cm(11.5),Cm(30),Cm(6),'📷 图 1-53 QPPB 实现原理示意图')

# S9 Key Points
s=prs.slides.add_slide(prs.slide_layouts[6]); TB(s,'2.2 关键技术点'); SN(s,9); FT(s)
kps=[('Behavior ID','不同的流动作对应不同的 Behavior ID，\n存储在 FIB 表项中。数据转发时根据目的网络\n从 FIB 获取 Behavior ID，执行相应流动作'),
     ('QoS Local-ID','QPPB 策略中绑定 qos-local-id 与 behavior，\n实现路由属性与 QoS 策略的关联'),
     ('策略传递机制','路由发送端：通过 route-policy 设置\nAS_PATH / Community / Ext-Community\n路由接收端：通过 route-policy import 匹配属性\napply qos-local-id')]
for i,(t,tx) in enumerate(kps):
    y=Cm(3+i*4.5); S(s,Cm(2),y,Cm(6),Cm(1.2),DB,t,16,WH,True)
    T(s,Cm(8.5),y+Cm(0.1),Cm(23),Cm(1.2),tx,12,BK,False)
T(s,Cm(2),Cm(15.5),Cm(30),Cm(2.5),
  '⚠️ 重要说明：QPPB技术实际并没有在 BGP 路由信息中发送 QoS 策略，\n只是在路由发送方通过对通告的路由设置路由属性，\n在路由接收方根据目的网段的路由属性设置 QoS 策略。',12,BK,False)

# S10 Up vs Down
s=prs.slides.add_slide(prs.slide_layouts[6]); TB(s,'2.3 上行 vs 下行 QPPB'); SN(s,10); FT(s)
td=[['方向','配置命令','查表依据','应用场景'],
    ['上行\n（inbound）','qppb-policy policy\nsource inbound','根据源 IP\n查路由表','用户→ISP\n流量计费'],
    ['下行\n（outbound）','qppb-policy policy\noutbound','根据目的 IP\n查路由表','ISP→用户\n流量计费'],
    ['基于 IP 优先级','qppb-policy ip-precedence\nsource','根据源/目的地址','按优先级分类']]
tbl=s.shapes.add_table(len(td),4,Cm(2),Cm(3),Cm(30),Cm(10)).table
tbl.columns[0].width=Cm(5); tbl.columns[1].width=Cm(10); tbl.columns[2].width=Cm(7); tbl.columns[3].width=Cm(8)
for ri,rd in enumerate(td):
    for ci,ct in enumerate(rd):
        c=tbl.cell(ri,ci); c.text=''; p=c.text_frame.paragraphs[0]; r=p.add_run(); sf(r,11,BK,False); r.text=ct
        p.alignment=PP_ALIGN.CENTER; c.vertical_anchor=MSO_ANCHOR.MIDDLE
        if ri==0: c.fill.solid(); c.fill.fore_color.rgb=DB
        elif ri%2==0: c.fill.solid(); c.fill.fore_color.rgb=LG

# S11 Part3 Title
s=prs.slides.add_slide(prs.slide_layouts[6])
bg=s.background; bg.fill.solid(); bg.fill.fore_color.rgb=DB
T(s,Cm(4),Cm(7),Cm(26),Cm(3),'第三部分',36,WH,True,PP_ALIGN.CENTER)
T(s,Cm(4),Cm(10.5),Cm(26),Cm(2),'QPPB 典型应用场景',24,LB,False,PP_ALIGN.CENTER); SN(s,11)

# S12 App1
s=prs.slides.add_slide(prs.slide_layouts[6]); TB(s,'3.1 典型应用一：AS 域间流量分类（图 1-54）'); SN(s,12); FT(s)
T(s,Cm(2),Cm(3),Cm(30),Cm(1.5),
  '场景：使用 QPPB 可以方便地在 AS100 的边缘设备对 AS 域间的流量进行流分类。\n例如要在 Node-C 上对 AS200 和 AS400 之间的流量进行限速。',14,BK,False)
T(s,Cm(2),Cm(5),Cm(30),Cm(1),'配置方案：',16,DB,True)
for i,tx in enumerate(['AS200 → AS400 方向：在 Node-C 上的 AS100 域内所有接口使能针对源地址的 QPPB',
                       'AS400 → AS200 方向：在 Node-C 上与 AS400 相连的接口使能针对目的地址的 QPPB']):
    S(s,Cm(2),Cm(6.5+i*1.5),Cm(30),Cm(1.2),LB,f'{i+1}. {tx}',12,DB,False,PP_ALIGN.LEFT)
T(s,Cm(2),Cm(10),Cm(30),Cm(1.5),
  '⚠️ 须知：查 FIB 转发的是针对上行流量而不是下行流量，\n因此使能 QPPB 的接口是流量上行的接口。',13,RD,True)
IP(s,Cm(2),Cm(12),Cm(30),Cm(5),'📷 图 1-54 AS 域间流量分类组网示意图')

# S13 App2
s=prs.slides.add_slide(prs.slide_layouts[6]); TB(s,'3.2 典型应用二：L3VPN 流量分类（图 1-55）'); SN(s,13); FT(s)
T(s,Cm(2),Cm(3),Cm(30),Cm(1.5),
  '场景：QPPB技术在 BGP/MPLS L3VPN 组网环境中的应用。\n当 PE 连接多个 VPN 时，可以对某个 VPN-instance 在路由发布时设置 Community 等属性。',14,BK,False)
T(s,Cm(2),Cm(5),Cm(30),Cm(1),'配置方案：',16,DB,True)
for i,tx in enumerate(['PE 连接多个 VPN 时，对某个 VPN-instance 在路由发布时设置 Community 等属性',
                       '远端 PE 接收到路由信息后将路由及 QoS 参数设置到 FIB 表项中',
                       '使得从 CE 来的流量在转发时能执行相应的 QoS 动作',
                       '不同的 VPN 可获得不同的服务质量']):
    S(s,Cm(2),Cm(6.5+i*1.5),Cm(30),Cm(1.2),LB,f'{i+1}. {tx}',12,DB,False,PP_ALIGN.LEFT)
IP(s,Cm(2),Cm(13),Cm(30),Cm(5),'📷 图 1-55 L3VPN 流量分类组网示意图')

# S14 App3
s=prs.slides.add_slide(prs.slide_layouts[6]); TB(s,'3.3 典型应用三：用户→ISP 的流量计费（图 1-56）'); SN(s,14); FT(s)
T(s,Cm(2),Cm(3),Cm(30),Cm(1),'场景：QPPB技术应用于用户到 ISP 的流量计费场景。',14,BK,False)
T(s,Cm(2),Cm(4.5),Cm(30),Cm(1),'配置方案：',16,DB,True)
for i,tx in enumerate(['通过 BGP 协议，发布路由时携带团体属性',
                       '引入 BGP 路由时，匹配团体属性，在路由表中设置 Behavior ID',
                       '配置 qppb-policy，匹配 qos-local-id，配置统计/CAR/Remark 等动作',
                       '在流量入口方向使能基于目的地址的 QPPB',
                       '在用户侧接口的 inbound 方向应用 qppb-policy']):
    S(s,Cm(2),Cm(6+i*1.3),Cm(30),Cm(1.1),LB,f'{i+1}. {tx}',11,DB,False,PP_ALIGN.LEFT)
T(s,Cm(2),Cm(13.5),Cm(30),Cm(2),
  '转发过程：目的 IP 查路由表 → 获取 Behavior ID → 匹配 qppb-policy → 执行统计/CAR/Remark 等动作',12,BK,False)
IP(s,Cm(2),Cm(15.5),Cm(30),Cm(3),'📷 图 1-56 用户→ISP 流量计费组网示意图')

# S15 App4
s=prs.slides.add_slide(prs.slide_layouts[6]); TB(s,'3.4 典型应用四：ISP→用户的流量计费（图 1-57）'); SN(s,15); FT(s)
T(s,Cm(2),Cm(3),Cm(30),Cm(1),'场景：QPPB技术应用于 ISP 到用户的流量计费场景。',14,BK,False)
T(s,Cm(2),Cm(4.5),Cm(30),Cm(1),'配置方案：',16,DB,True)
for i,tx in enumerate(['通过 BGP 协议，发布路由时携带团体属性',
                       '引入 BGP 路由时，匹配团体属性，在路由表中设置 Behavior ID',
                       '配置 qppb-policy，匹配 qos-local-id，配置统计/CAR/Remark 等动作',
                       '在流量入接口方向使能基于源地址的 QPPB',
                       '在用户侧接口的 outbound 方向应用 qppb-policy']):
    S(s,Cm(2),Cm(6+i*1.3),Cm(30),Cm(1.1),LB,f'{i+1}. {tx}',11,DB,False,PP_ALIGN.LEFT)
T(s,Cm(2),Cm(13.5),Cm(30),Cm(2),
  '转发过程：源 IP 查路由表 → 获取 Behavior ID → 内部交换传递到出接口 → 匹配 qppb-policy → 执行动作',12,BK,False)
IP(s,Cm(2),Cm(15.5),Cm(30),Cm(3),'📷 图 1-57 ISP→用户流量计费组网示意图')

out='/home/admin/.openclaw/workspace/QPPB技术详解.pptx'
prs.save(out)
print(f'✅ PPT 已生成: {out}')
