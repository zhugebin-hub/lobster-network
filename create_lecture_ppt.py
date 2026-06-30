#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
生成"从小龙虾到数字员工——AI智能体在商务管理与跨境电商中的应用实践"讲座 PPT
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ============ 配色方案 ============
PRIMARY = RGBColor(0xE8, 0x6B, 0x2D)    # 小龙虾橙
DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)     # 深蓝黑
ACCENT = RGBColor(0xFF, 0xA5, 0x00)      # 金色
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF0, 0xF0, 0xF5)
TEXT_DARK = RGBColor(0x2C, 0x2C, 0x3E)
TEXT_LIGHT = RGBColor(0x55, 0x55, 0x70)
CARD_BG = RGBColor(0xFF, 0xF8, 0xF0)
GREEN = RGBColor(0x27, 0xAE, 0x60)
BLUE = RGBColor(0x29, 0x80, 0xB9)
RED = RGBColor(0xE7, 0x4C, 0x3C)

def add_bg(slide, color=DARK_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape_bg(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name='微软雅黑'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_multi_text(slide, left, top, width, height, lines, font_size=16, color=WHITE, line_spacing=1.5, font_name='微软雅黑'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = Pt(font_size * (line_spacing - 1))
    return txBox

def add_card(slide, left, top, width, height, title, content_lines, title_color=PRIMARY, bg_color=CARD_BG):
    shape = add_shape_bg(slide, left, top, width, height, bg_color)
    add_text_box(slide, left + Inches(0.2), top + Inches(0.15), width - Inches(0.4), Inches(0.5),
                 title, font_size=16, color=title_color, bold=True)
    add_multi_text(slide, left + Inches(0.2), top + Inches(0.6), width - Inches(0.4), height - Inches(0.8),
                   content_lines, font_size=13, color=TEXT_DARK, line_spacing=1.4)
    return shape

# ============ Slide 1: 封面 ============
slide1 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide1, DARK_BG)
add_shape_bg(slide1, Inches(0), Inches(0), Inches(13.333), Inches(0.08), PRIMARY)
add_shape_bg(slide1, Inches(0), Inches(7.42), Inches(13.333), Inches(0.08), PRIMARY)
add_text_box(slide1, Inches(5.5), Inches(1.0), Inches(2.5), Inches(1.5), "\U0001F99E", font_size=72, alignment=PP_ALIGN.CENTER)
add_text_box(slide1, Inches(1.5), Inches(2.5), Inches(10.3), Inches(1.2),
             "从小龙虾到数字员工", font_size=44, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide1, Inches(1.5), Inches(3.6), Inches(10.3), Inches(1.0),
             "AI智能体在商务管理与跨境电商中的应用实践", font_size=24, color=ACCENT, alignment=PP_ALIGN.CENTER)
add_shape_bg(slide1, Inches(4.5), Inches(4.5), Inches(4.3), Inches(0.04), PRIMARY)
add_text_box(slide1, Inches(3), Inches(4.8), Inches(7.3), Inches(0.5),
             "浙江越秀外国语学院", font_size=20, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
add_text_box(slide1, Inches(3), Inches(5.4), Inches(7.3), Inches(0.5),
             "商务管理 + 跨境电商 | 约100人", font_size=16, color=TEXT_LIGHT, alignment=PP_ALIGN.CENTER)
add_text_box(slide1, Inches(3), Inches(5.9), Inches(7.3), Inches(0.5),
             "2026年6月2日（周二）", font_size=16, color=TEXT_LIGHT, alignment=PP_ALIGN.CENTER)
add_text_box(slide1, Inches(3), Inches(6.6), Inches(7.3), Inches(0.5),
             "主讲人：诸葛斌", font_size=14, color=TEXT_LIGHT, alignment=PP_ALIGN.CENTER)

# ============ Slide 2: 目录 ============
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide2, DARK_BG)
add_shape_bg(slide2, Inches(0), Inches(0), Inches(13.333), Inches(1.2), PRIMARY)
add_text_box(slide2, Inches(1), Inches(0.3), Inches(5), Inches(0.8), "\U0001F4CB 讲座目录", font_size=32, color=WHITE, bold=True)

toc_items = [
    ("01", "AI 时代：从工具到智能体的范式迁移"),
    ("02", "小龙虾智能体：什么是数字员工"),
    ("03", "璀璨臻选：AI 珠宝定制平台案例"),
    ("04", "数字员工在商务管理中的应用"),
    ("05", "数字员工在跨境电商中的应用"),
    ("06", "小龙虾三部曲：人机协作方法论"),
    ("07", "AI 时代大学生的职业新定位"),
    ("08", "互动体验与 Q&A"),
]

for i, (num, title) in enumerate(toc_items):
    y = Inches(1.5) + Inches(i * 0.7)
    add_text_box(slide2, Inches(2), y, Inches(0.8), Inches(0.5), num, font_size=24, color=PRIMARY, bold=True)
    add_text_box(slide2, Inches(3), y, Inches(8), Inches(0.5), title, font_size=20, color=WHITE)
    if i < len(toc_items) - 1:
        add_shape_bg(slide2, Inches(2), y + Inches(0.55), Inches(9), Inches(0.01), RGBColor(0x33, 0x33, 0x55))

# ============ Slide 3: AI 时代背景 ============
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide3, DARK_BG)
add_shape_bg(slide3, Inches(0), Inches(0), Inches(13.333), Inches(1.2), PRIMARY)
add_text_box(slide3, Inches(1), Inches(0.3), Inches(10), Inches(0.8), "\U0001F30A AI 时代：从工具到智能体的范式迁移", font_size=30, color=WHITE, bold=True)

phases = [
    ("2022", "对话框时代\nChatGPT 引爆", "问答交互"),
    ("2024", "大模型普及\n多模态崛起", "内容生成"),
    ("2025", "智能体元年\nAgent 爆发", "自主决策"),
    ("2026", "数字员工\n大规模部署", "人机协作"),
    ("2027", "智能体普及率\n超 70%", "全面融合"),
]

for i, (year, desc, tag) in enumerate(phases):
    x = Inches(0.8) + Inches(i * 2.4)
    add_text_box(slide3, x, Inches(1.8), Inches(2), Inches(0.5), year, font_size=22, color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide3, x, Inches(2.3), Inches(2.2), Inches(1.2), desc, font_size=14, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_shape_bg(slide3, x + Inches(0.2), Inches(3.5), Inches(1.6), Inches(0.45), PRIMARY)
    add_text_box(slide3, x + Inches(0.2), Inches(3.5), Inches(1.6), Inches(0.45), tag, font_size=12, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    if i < len(phases) - 1:
        add_shape_bg(slide3, x + Inches(2.0), Inches(2.1), Inches(0.4), Inches(0.04), ACCENT)

add_shape_bg(slide3, Inches(0.8), Inches(4.5), Inches(11.7), Inches(2.5), RGBColor(0x1E, 0x1E, 0x3A))
add_text_box(slide3, Inches(1.2), Inches(4.7), Inches(11), Inches(0.5), "\U0001F4CB 国家战略：智能体普及年", font_size=22, color=ACCENT, bold=True)

policy_lines = [
    "\u2022 国发〔2025〕11号文：到2027年智能体应用普及率超70%",
    "\u2022 2026年政府工作报告首次提出【打造智能经济新形态】",
    "\u2022 工信部【模数共振】行动：培育垂直领域大模型和专用智能体",
    "\u2022 算力纳入国家【六网】新基建，算力即水电",
]
add_multi_text(slide3, Inches(1.2), Inches(5.3), Inches(11), Inches(1.5), policy_lines, font_size=15, color=LIGHT_GRAY, line_spacing=1.6)

# ============ Slide 4: 什么是小龙虾智能体 ============
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide4, DARK_BG)
add_shape_bg(slide4, Inches(0), Inches(0), Inches(13.333), Inches(1.2), PRIMARY)
add_text_box(slide4, Inches(1), Inches(0.3), Inches(10), Inches(0.8), "\U0001F99E 小龙虾智能体：什么是数字员工", font_size=30, color=WHITE, bold=True)

add_card(slide4, Inches(0.8), Inches(1.5), Inches(5.5), Inches(3.0),
         "\U0001F914 什么是数字员工？",
         ["数字员工 = AI智能体 + 业务能力",
          "不是简单的聊天机器人，而是能理解、",
          "决策、执行的【AI 同事】",
          "24 小时在线，不知疲倦",
          "能帮你开会、写报告、查数据、管项目"],
         bg_color=RGBColor(0x1E, 0x1E, 0x3A))

add_card(slide4, Inches(6.8), Inches(1.5), Inches(5.7), Inches(3.0),
         "\U0001F4CA 传统 AI vs 数字员工",
         ["传统 AI：问答工具，用完即走",
          "数字员工：持续在岗，主动工作",
          "",
          "传统 AI：通用能力，不懂业务",
          "数字员工：懂业务、懂流程、懂权限"],
         bg_color=RGBColor(0x1E, 0x1E, 0x3A))

add_shape_bg(slide4, Inches(0.8), Inches(4.8), Inches(11.7), Inches(2.3), RGBColor(0x1E, 0x1E, 0x3A))
add_text_box(slide4, Inches(1.2), Inches(4.9), Inches(10), Inches(0.5), "\U0001F527 小龙虾智能体核心能力", font_size=20, color=ACCENT, bold=True)

caps = [
    ("智能对话", "通义千问企业版\n专业术语理解"),
    ("工作流编排", "低代码可视化\n业务人员可配置"),
    ("知识库 RAG", "企业文档向量化\n准确检索90%+"),
    ("API 连接器", "钉钉/阿里云\n原生无缝对接"),
    ("权限引擎", "细粒度RBAC\n完整审计日志"),
]

for i, (title, desc) in enumerate(caps):
    x = Inches(1.0) + Inches(i * 2.2)
    add_card(slide4, x, Inches(5.4), Inches(2.0), Inches(1.5), title, desc.split('\n'),
             bg_color=RGBColor(0x25, 0x25, 0x45))

# ============ Slide 5: 璀璨臻选项目介绍 ============
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide5, DARK_BG)
add_shape_bg(slide5, Inches(0), Inches(0), Inches(13.333), Inches(1.2), PRIMARY)
add_text_box(slide5, Inches(1), Inches(0.3), Inches(10), Inches(0.8), "\U0001F48E 璀璨臻选：AI 高级珠宝定制平台", font_size=30, color=WHITE, bold=True)

add_card(slide5, Inches(0.8), Inches(1.5), Inches(5.5), Inches(2.5),
         "\U0001F3AF 项目定位",
         ["重构义乌小商品产业链的 AI 科技项目",
          "以 LLM + AIGC 视觉生成技术为核心",
          "首创【AI 设计平台 + OpenClaw 数字员工派遣】",
          "双轮驱动模式",
          "7 项已受理发明专利 + 1 项已筹稿"],
         bg_color=RGBColor(0x1E, 0x1E, 0x3A))

add_card(slide5, Inches(6.8), Inches(1.5), Inches(5.7), Inches(2.5),
         "\U0001F504 双轮驱动模式",
         ["C 端消费者：【一句话生成专属珠宝】",
          "  -> 普惠定制体验，零设计门槛",
          "",
          "B 端义乌中小微企业：云上小龙虾数字员工",
          "  -> 市场洞察 + 智能跟单 + 柔性生产"],
         bg_color=RGBColor(0x1E, 0x1E, 0x3A))

add_shape_bg(slide5, Inches(0.8), Inches(4.3), Inches(11.7), Inches(2.8), RGBColor(0x1E, 0x1E, 0x3A))
add_text_box(slide5, Inches(1.2), Inches(4.4), Inches(10), Inches(0.5), "\U0001F3AF 行业痛点 vs 解决方案", font_size=20, color=ACCENT, bold=True)

pain_points = [
    ("C 端：设计壁垒高", "Z世代超78%对定制感兴趣\n但传统定制门槛高、周期长、价格贵"),
    ("B 端：低价内卷", "义乌饰品产业依赖薄利多销\n缺乏数字化转型能力"),
    ("生产端：柔性生产难", "小批量、个性化需求\n传统产线无法快速响应"),
]

for i, (title, desc) in enumerate(pain_points):
    x = Inches(1.0) + Inches(i * 3.8)
    add_card(slide5, x, Inches(5.0), Inches(3.5), Inches(1.8), title, desc.split('\n'),
             title_color=RED, bg_color=RGBColor(0x25, 0x25, 0x45))

# ============ Slide 6: 数字员工在商务管理中的应用 ============
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide6, DARK_BG)
add_shape_bg(slide6, Inches(0), Inches(0), Inches(13.333), Inches(1.2), PRIMARY)
add_text_box(slide6, Inches(1), Inches(0.3), Inches(10), Inches(0.8), "\U0001F4BC 数字员工在商务管理中的应用", font_size=30, color=WHITE, bold=True)

scenarios = [
    ("会议管理", "录音转写 -> AI提取决议\n-> 创建钉钉任务 -> 自动提醒", "纪要整理\n93%\u2193"),
    ("项目管理", "低代码配置预警规则\n-> Teambition进度同步\n-> 周报自动推送", "规则配置\n90%\u2193"),
    ("团队管理", "组织架构自动同步\n-> 生日个性化祝福\n-> 绩效数据聚合", "信息维护\n100%\u2193"),
    ("业务运营", "DataWorks抽取数据\n-> QuickBI分析\n-> AI生成日报解读", "日报编制\n92%\u2193"),
    ("知识管理", "文档自动切片向量化\n-> 语义检索 + 答案生成\n-> 来源可追溯", "制度查询\n95%\u2193"),
]

for i, (title, desc, result) in enumerate(scenarios):
    x = Inches(0.5) + Inches(i * 2.5)
    add_card(slide6, x, Inches(1.5), Inches(2.3), Inches(3.5), title, desc.split('\n'),
             bg_color=RGBColor(0x1E, 0x1E, 0x3A))
    add_shape_bg(slide6, x + Inches(0.3), Inches(5.1), Inches(1.7), Inches(0.5), GREEN)
    add_text_box(slide6, x + Inches(0.3), Inches(5.1), Inches(1.7), Inches(0.5), result, font_size=13, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_shape_bg(slide6, Inches(0.8), Inches(5.8), Inches(11.7), Inches(1.3), RGBColor(0x1E, 0x1E, 0x3A))
add_text_box(slide6, Inches(1.2), Inches(5.9), Inches(10), Inches(0.5), "\U0001F4A1 核心价值：让数字员工成为【24小时不打烊】的商务助手", font_size=18, color=ACCENT, bold=True)
add_text_box(slide6, Inches(1.2), Inches(6.4), Inches(11), Inches(0.6),
             "从【工具辅助】到【智能代理】——商务管理的范式迁移", font_size=15, color=LIGHT_GRAY)

# ============ Slide 7: 数字员工在跨境电商中的应用 ============
slide7 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide7, DARK_BG)
add_shape_bg(slide7, Inches(0), Inches(0), Inches(13.333), Inches(1.2), PRIMARY)
add_text_box(slide7, Inches(1), Inches(0.3), Inches(10), Inches(0.8), "\U0001F30D 数字员工在跨境电商中的应用", font_size=30, color=WHITE, bold=True)

cross_border = [
    ("智能选品", "AI分析市场趋势\n竞品监控\n爆款预测", "选品效率\n提升5倍"),
    ("AI文案", "多语言产品描述\n社交媒体内容\n广告文案生成", "内容生产\n成本降80%"),
    ("智能客服", "24小时多语言客服\n订单查询\n售后处理", "响应时间\n<30秒"),
    ("数据分析", "销售数据实时分析\n库存预警\n利润核算", "决策效率\n提升3倍"),
    ("自动化运营", "上架/下架管理\n价格优化\n促销活动", "运营效率\n提升10倍"),
]

for i, (title, desc, result) in enumerate(cross_border):
    x = Inches(0.5) + Inches(i * 2.5)
    add_card(slide7, x, Inches(1.5), Inches(2.3), Inches(3.5), title, desc.split('\n'),
             bg_color=RGBColor(0x1E, 0x1E, 0x3A))
    add_shape_bg(slide7, x + Inches(0.3), Inches(5.1), Inches(1.7), Inches(0.5), BLUE)
    add_text_box(slide7, x + Inches(0.3), Inches(5.1), Inches(1.7), Inches(0.5), result, font_size=13, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_shape_bg(slide7, Inches(0.8), Inches(5.8), Inches(11.7), Inches(1.3), RGBColor(0x1E, 0x1E, 0x3A))
add_text_box(slide7, Inches(1.2), Inches(5.9), Inches(10), Inches(0.5), "\U0001F31F 璀璨臻选案例：AI 珠宝定制的跨境电商实践", font_size=18, color=ACCENT, bold=True)
add_text_box(slide7, Inches(1.2), Inches(6.4), Inches(11), Inches(0.6),
             "一句话生成专属珠宝设计 -> AI 视觉生成 -> 义乌柔性生产 -> 全球发货", font_size=15, color=LIGHT_GRAY)

# ============ Slide 8: 小龙虾三部曲方法论 ============
slide8 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide8, DARK_BG)
add_shape_bg(slide8, Inches(0), Inches(0), Inches(13.333), Inches(1.2), PRIMARY)
add_text_box(slide8, Inches(1), Inches(0.3), Inches(10), Inches(0.8), "\U0001F4C8 小龙虾三部曲：人机协作方法论", font_size=30, color=WHITE, bold=True)

steps = [
    ("第一部", "认 知", "什么是小龙虾\n能力边界在哪\n数据安全与隐私", "准确判断\n哪些任务适合\n交给小龙虾"),
    ("第二部", "技 能", "有效提问技巧\n多轮迭代优化\n输出验证方法", "提问效率\n提升50%\n独立配置\n简单工作流"),
    ("第三部", "理 念", "人机协作共生\n双向成长\n最终独立", "形成正确的\n人机协作观\n主动优化\n协作模式"),
]

for i, (step, name, content, outcome) in enumerate(steps):
    x = Inches(0.8) + Inches(i * 4.0)
    add_shape_bg(slide8, x, Inches(1.5), Inches(3.6), Inches(0.5), PRIMARY)
    add_text_box(slide8, x, Inches(1.5), Inches(3.6), Inches(0.5), step, font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_shape_bg(slide8, x, Inches(2.0), Inches(3.6), Inches(0.8), RGBColor(0x25, 0x25, 0x45))
    add_text_box(slide8, x, Inches(2.0), Inches(3.6), Inches(0.8), name, font_size=28, color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER)
    add_card(slide8, x, Inches(2.9), Inches(3.6), Inches(2.0), "学什么", content.split('\n'),
             bg_color=RGBColor(0x1E, 0x1E, 0x3A))
    add_card(slide8, x, Inches(5.0), Inches(3.6), Inches(1.8), "\U0001F3AF 产出", outcome.split('\n'),
             title_color=GREEN, bg_color=RGBColor(0x1E, 0x1E, 0x3A))
    if i < 2:
        add_text_box(slide8, x + Inches(3.6), Inches(2.2), Inches(0.4), Inches(0.6), "\u2192", font_size=36, color=ACCENT, alignment=PP_ALIGN.CENTER)

# ============ Slide 9: AI时代大学生职业新定位 ============
slide9 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide9, DARK_BG)
add_shape_bg(slide9, Inches(0), Inches(0), Inches(13.333), Inches(1.2), PRIMARY)
add_text_box(slide9, Inches(1), Inches(0.3), Inches(10), Inches(0.8), "\U0001F393 AI 时代大学生的职业新定位", font_size=30, color=WHITE, bold=True)

add_card(slide9, Inches(0.8), Inches(1.5), Inches(5.5), Inches(3.0),
         "\u26A0\uFE0F AI 时代带来的挑战",
         ["\u2022 传统岗位被 AI 替代（客服、文案、数据分析）",
          "\u2022 企业对【AI 素养】的要求越来越高",
          "\u2022 只会【用工具】不够，要会【指挥工具】",
          "\u2022 人机协作能力成为核心竞争力"],
         title_color=RED, bg_color=RGBColor(0x1E, 0x1E, 0x3A))

add_card(slide9, Inches(6.8), Inches(1.5), Inches(5.7), Inches(3.0),
         "\U0001F31F AI 时代带来的机遇",
         ["\u2022 新职业涌现：AI 训练师、智能体工程师、提示词工程师",
          "\u2022 一人公司成为可能：AI 帮你做客服、运营、设计",
          "\u2022 创业门槛降低：AI 弥补团队能力短板",
          "\u2022 效率革命：1个人 = 1个团队"],
         title_color=GREEN, bg_color=RGBColor(0x1E, 0x1E, 0x3A))

add_shape_bg(slide9, Inches(0.8), Inches(4.8), Inches(11.7), Inches(2.3), RGBColor(0x1E, 0x1E, 0x3A))
add_text_box(slide9, Inches(1.2), Inches(4.9), Inches(10), Inches(0.5), "\U0001F511 未来人才能力模型", font_size=22, color=ACCENT, bold=True)

skills = [
    ("AI 素养", "理解 AI 能力边界\n掌握有效提问技巧"),
    ("业务洞察", "理解行业痛点\n发现 AI 应用场景"),
    ("人机协作", "与 AI 高效配合\n形成工作闭环"),
    ("创新思维", "用 AI 创造新价值\n而非替代旧流程"),
    ("数据思维", "用数据驱动决策\n用 AI 放大洞察"),
]

for i, (title, desc) in enumerate(skills):
    x = Inches(0.8) + Inches(i * 2.4)
    add_card(slide9, x, Inches(5.4), Inches(2.2), Inches(1.5), title, desc.split('\n'),
             bg_color=RGBColor(0x25, 0x25, 0x45))

# ============ Slide 10: 互动体验 ============
slide10 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide10, DARK_BG)
add_shape_bg(slide10, Inches(0), Inches(0), Inches(13.333), Inches(1.2), PRIMARY)
add_text_box(slide10, Inches(1), Inches(0.3), Inches(10), Inches(0.8), "\U0001F3AE 互动体验与 Q&A", font_size=30, color=WHITE, bold=True)

add_card(slide10, Inches(1.0), Inches(1.5), Inches(5.5), Inches(2.5),
         "\U0001F99E 现场体验：一句话定制珠宝",
         ["1. 说出你的设计理念（文字描述）",
          "2. AI 实时生成珠宝设计图",
          "3. 查看生产参数与成本估算",
          "4. 体验数字员工全流程服务"],
         bg_color=RGBColor(0x1E, 0x1E, 0x3A))

add_card(slide10, Inches(7.0), Inches(1.5), Inches(5.3), Inches(2.5),
         "\u2753 自由提问",
         ["关于 AI智能体、数字员工、",
          "璀璨臻选项目、职业发展",
          "任何问题，欢迎提问！",
          "",
          "\U0001F381 提问同学有机会获得",
          "   AI 生成的专属珠宝设计图"],
         bg_color=RGBColor(0x1E, 0x1E, 0x3A))

add_shape_bg(slide10, Inches(1.0), Inches(4.3), Inches(11.3), Inches(2.5), RGBColor(0x1E, 0x1E, 0x3A))
add_text_box(slide10, Inches(1.5), Inches(4.5), Inches(10), Inches(0.5), "\U0001F4EC 获取更多学习资料", font_size=20, color=ACCENT, bold=True)

contact_lines = [
    "\u2022 小龙虾智能体体验平台：阿里云百炼平台",
    "\u2022 璀璨臻选项目演示：现场扫码体验",
    "\u2022 AI 学习资源推荐：讲座后发送资料包",
    "\u2022 交流群：扫码加入【AI智能体学习群】",
]
add_multi_text(slide10, Inches(1.5), Inches(5.1), Inches(10), Inches(1.5), contact_lines, font_size=16, color=LIGHT_GRAY, line_spacing=1.6)

# ============ Slide 11: 致谢 ============
slide11 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide11, DARK_BG)
add_shape_bg(slide11, Inches(0), Inches(0), Inches(13.333), Inches(0.08), PRIMARY)
add_shape_bg(slide11, Inches(0), Inches(7.42), Inches(13.333), Inches(0.08), PRIMARY)
add_text_box(slide11, Inches(5.5), Inches(1.5), Inches(2.5), Inches(1.5), "\U0001F99E", font_size=72, alignment=PP_ALIGN.CENTER)
add_text_box(slide11, Inches(2), Inches(3.0), Inches(9.3), Inches(1.0),
             "谢谢聆听！", font_size=48, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_shape_bg(slide11, Inches(4.5), Inches(4.0), Inches(4.3), Inches(0.04), PRIMARY)
add_text_box(slide11, Inches(2), Inches(4.3), Inches(9.3), Inches(0.8),
             "让小龙虾成为你的得力数字员工", font_size=24, color=ACCENT, alignment=PP_ALIGN.CENTER)
add_text_box(slide11, Inches(2), Inches(5.5), Inches(9.3), Inches(0.5),
             "浙江越秀外国语学院 | 商务管理 + 跨境电商", font_size=16, color=TEXT_LIGHT, alignment=PP_ALIGN.CENTER)
add_text_box(slide11, Inches(2), Inches(6.0), Inches(9.3), Inches(0.5),
             "2026年6月2日 | 主讲人：诸葛斌", font_size=16, color=TEXT_LIGHT, alignment=PP_ALIGN.CENTER)

# ============ 保存 ============
output_path = "/home/admin/.openclaw/workspace/从小龙虾到数字员工_AI智能体讲座.pptx"
prs.save(output_path)
print(f"PPT 已保存: {output_path}")
