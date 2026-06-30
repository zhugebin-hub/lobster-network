#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""基于模板生成 PPT"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from lxml import etree

TEMPLATE = "/home/admin/.openclaw/media/inbound/3e441551-6f97-462e-b2ba-dd1e57f3756d.pptx"
OUT = "/home/admin/.openclaw/workspace/科学衔接赋能成长.pptx"

prs = Presentation(TEMPLATE)
LAYOUTS = {l.name: i for i, l in enumerate(prs.slide_layouts)}
print("布局:", list(LAYOUTS.keys()))

def set_text(shape, text, size=18, color=None, bold=None, align=None):
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    if color: p.font.color.rgb = color
    if bold is not None: p.font.bold = bold
    if align: p.alignment = align
    p.font.name = '微软雅黑'
    return p

def add_para(tf, text, size=18, color=None, bold=None, align=None, space_before=4):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    if color: p.font.color.rgb = color
    if bold is not None: p.font.bold = bold
    if align: p.alignment = align
    p.font.name = '微软雅黑'
    p.space_before = Pt(space_before)
    return p

# ═══════════════════════════════════════════
# Slide 1: 封面页
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[LAYOUTS['封面页-1_4']])
for shape in slide.placeholders:
    ph = shape.placeholder_format.idx
    if ph == 14: set_text(shape, "科学衔接  赋能成长", size=40, bold=True, align=PP_ALIGN.CENTER)
    elif ph == 15: set_text(shape, "——让孩子有准备地进入小学", size=24, align=PP_ALIGN.CENTER)
    elif ph == 16: set_text(shape, "2026.05", size=18, align=PP_ALIGN.CENTER)
    elif ph == 17: set_text(shape, "高照实验学校  周  高", size=18, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════
# Slide 2: 章节页 - 您的焦虑
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[LAYOUTS['章节页-1_4']])
for shape in slide.placeholders:
    if shape.placeholder_format.idx == 14: set_text(shape, "01", size=48, bold=True)
    elif shape.placeholder_format.idx == 15: set_text(shape, "您的焦虑，我们都懂", size=32, bold=True)

# ═══════════════════════════════════════════
# Slide 3: 三个焦虑问题
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[LAYOUTS['标题和内容']])
for shape in slide.placeholders:
    if shape.placeholder_format.idx == 0:
        set_text(shape, "家长的三大焦虑", size=32, bold=True)
    elif shape.placeholder_format.idx == 1:
        tf = shape.text_frame; tf.clear()
        add_para(tf, "📖  要不要提前学拼音？", size=24, bold=True, space_before=16)
        add_para(tf, "😟  坐不住被老师批评怎么办？", size=24, bold=True, space_before=16)
        add_para(tf, "😰  跟不上会自卑吗？", size=24, bold=True, space_before=16)
        add_para(tf, "", size=12, space_before=16)
        add_para(tf, "这些焦虑，是爱的开始。", size=20, color=RGBColor(0x43,0xA0,0x47), bold=True, space_before=8)

# ═══════════════════════════════════════════
# Slide 4: 章节页 - 核心观点
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[LAYOUTS['章节页-1_4']])
for shape in slide.placeholders:
    if shape.placeholder_format.idx == 14: set_text(shape, "02", size=48, bold=True)
    elif shape.placeholder_format.idx == 15: set_text(shape, "一个核心观点", size=32, bold=True)

# ═══════════════════════════════════════════
# Slide 5: 核心观点
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[LAYOUTS['标题和内容']])
for shape in slide.placeholders:
    if shape.placeholder_format.idx == 0:
        set_text(shape, "幼小衔接 ≠ 知识抢跑", size=36, bold=True)
    elif shape.placeholder_format.idx == 1:
        tf = shape.text_frame; tf.clear()
        for title, desc in [("心理准备","建立上学的期待与安全感"),("习惯准备","作息规律、生活自理"),("能力准备","专注力、表达力、动手力"),("社会适应","规则意识、交往能力")]:
            add_para(tf, f"▎{title}", size=22, bold=True, space_before=12)
            add_para(tf, f"    {desc}", size=18, space_before=2)

# ═══════════════════════════════════════════
# Slide 6: 章节页 - 幼儿园 vs 小学
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[LAYOUTS['章节页-1_4']])
for shape in slide.placeholders:
    if shape.placeholder_format.idx == 14: set_text(shape, "03", size=48, bold=True)
    elif shape.placeholder_format.idx == 15: set_text(shape, "幼儿园 vs 小学：六个不一样", size=32, bold=True)

# ═══════════════════════════════════════════
# Slide 7: 比较 - 上
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[LAYOUTS['比较']])
for shape in slide.placeholders:
    ph = shape.placeholder_format.idx
    if ph == 0: set_text(shape, "幼儿园", size=28, bold=True)
    elif ph == 1:
        tf = shape.text_frame; tf.clear()
        add_para(tf, "🏫  学习环境", size=20, bold=True, space_before=10)
        add_para(tf, "游戏化、生活化", size=18, space_before=4)
        add_para(tf, "👩‍🏫  师生关系", size=20, bold=True, space_before=10)
        add_para(tf, "两教一保，细致照顾", size=18, space_before=4)
        add_para(tf, "⏰  作息时间", size=20, bold=True, space_before=10)
        add_para(tf, "午睡长，弹性", size=18, space_before=4)
    elif ph == 4: set_text(shape, "小学", size=28, bold=True)
    elif ph == 5:
        tf = shape.text_frame; tf.clear()
        add_para(tf, "🏫  学习环境", size=20, bold=True, space_before=10)
        add_para(tf, "班级授课、有纪律", size=18, space_before=4)
        add_para(tf, "👩‍🏫  师生关系", size=20, bold=True, space_before=10)
        add_para(tf, "班主任+科任，更自主", size=18, space_before=4)
        add_para(tf, "⏰  作息时间", size=20, bold=True, space_before=10)
        add_para(tf, "课时固定，午睡短", size=18, space_before=4)

# ═══════════════════════════════════════════
# Slide 8: 比较 - 下
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[LAYOUTS['比较']])
for shape in slide.placeholders:
    ph = shape.placeholder_format.idx
    if ph == 0: set_text(shape, "幼儿园", size=28, bold=True)
    elif ph == 1:
        tf = shape.text_frame; tf.clear()
        add_para(tf, "📚  学习方式", size=20, bold=True, space_before=10)
        add_para(tf, "玩中学", size=18, space_before=4)
        add_para(tf, "📏  行为规范", size=20, bold=True, space_before=10)
        add_para(tf, "相对宽松", size=18, space_before=4)
        add_para(tf, "⭐  评价方式", size=20, bold=True, space_before=10)
        add_para(tf, "过程性", size=18, space_before=4)
    elif ph == 4: set_text(shape, "小学", size=28, bold=True)
    elif ph == 5:
        tf = shape.text_frame; tf.clear()
        add_para(tf, "📚  学习方式", size=20, bold=True, space_before=10)
        add_para(tf, "听讲、作业、纸笔练习", size=18, space_before=4)
        add_para(tf, "📏  行为规范", size=20, bold=True, space_before=10)
        add_para(tf, "纪律、铃声、排队、值日", size=18, space_before=4)
        add_para(tf, "⭐  评价方式", size=20, bold=True, space_before=10)
        add_para(tf, "分数、考试、横向比较", size=18, space_before=4)

txBox = slide.shapes.add_textbox(Inches(1), Inches(6.2), Inches(11.3), Inches(0.6))
tf = txBox.text_frame
p = tf.paragraphs[0]; p.text = "这些是台阶，不是问题。"
p.font.size = Pt(20); p.font.color.rgb = RGBColor(0x43,0xA0,0x47)
p.font.bold = True; p.font.name = '微软雅黑'; p.alignment = PP_ALIGN.CENTER

# ═══════════════════════════════════════════
# Slide 9: 章节页 - 科学准备
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[LAYOUTS['章节页-1_4']])
for shape in slide.placeholders:
    if shape.placeholder_format.idx == 14: set_text(shape, "04", size=48, bold=True)
    elif shape.placeholder_format.idx == 15: set_text(shape, "科学准备：四个维度", size=32, bold=True)

# ═══════════════════════════════════════════
# Slide 10: 两栏 - 身心 + 生活
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[LAYOUTS['两栏内容']])
for shape in slide.placeholders:
    ph = shape.placeholder_format.idx
    if ph == 0: set_text(shape, "🏃 身心准备", size=28, bold=True)
    elif ph == 1:
        tf = shape.text_frame; tf.clear()
        add_para(tf, "愿意上学，情绪稳定", size=18, bold=True, space_before=4)
        add_para(tf, "• 用积极语言描述小学", size=18, space_before=8)
        add_para(tf, "• 带孩子熟悉上下学路线", size=18, space_before=8)
        add_para(tf, "• 告诉孩子成为小学生是骄傲的事", size=18, space_before=8)
        add_para(tf, "• 避免吓唬孩子", size=18, space_before=8)
    elif ph == 2: set_text(shape, "🎒 生活准备", size=28, bold=True)
    elif ph == 3:
        tf = shape.text_frame; tf.clear()
        add_para(tf, "自己的事情自己做", size=18, bold=True, space_before=4)
        add_para(tf, "• 独立如厕、整理书包、保管文具", size=18, space_before=8)
        add_para(tf, "• 提前2个月调整作息：21:00前入睡", size=18, space_before=8)
        add_para(tf, "• 简单劳动：擦桌、扫地、分碗筷", size=18, space_before=8)

# ═══════════════════════════════════════════
# Slide 11: 两栏 - 社会 + 学习
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[LAYOUTS['两栏内容']])
for shape in slide.placeholders:
    ph = shape.placeholder_format.idx
    if ph == 0: set_text(shape, "🤝 社会准备", size=28, bold=True)
    elif ph == 1:
        tf = shape.text_frame; tf.clear()
        add_para(tf, "会听·会说·会合作", size=18, bold=True, space_before=4)
        add_para(tf, "• 会听：指令只讲一遍，不重复", size=18, space_before=8)
        add_para(tf, "• 会说：老师我需要帮助", size=18, space_before=8)
        add_para(tf, "• 会合作：使用请谢谢对不起", size=18, space_before=8)
        add_para(tf, "• 冲突是学习社交的机会", size=18, space_before=8)
    elif ph == 2: set_text(shape, "📖 学习准备", size=28, bold=True)
    elif ph == 3:
        tf = shape.text_frame; tf.clear()
        add_para(tf, "兴趣比知识更重要", size=18, bold=True, space_before=4)
        add_para(tf, "• 不推荐：提前大量学拼音写字", size=18, space_before=8)
        add_para(tf, "• 推荐：亲子阅读、生活认字、数字游戏", size=18, space_before=8)
        add_para(tf, "• 手部锻炼：串珠 / 捏彩泥", size=18, space_before=8)
        add_para(tf, "• 牢记：纠正比新教更难", size=18, space_before=8)

# ═══════════════════════════════════════════
# Slide 12: 章节页 - 家长守则
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[LAYOUTS['章节页-1_4']])
for shape in slide.placeholders:
    if shape.placeholder_format.idx == 14: set_text(shape, "05", size=48, bold=True)
    elif shape.placeholder_format.idx == 15: set_text(shape, "家长守则", size=32, bold=True)

# ═══════════════════════════════════════════
# Slide 13: 三不要
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[LAYOUTS['标题和内容']])
for shape in slide.placeholders:
    if shape.placeholder_format.idx == 0:
        set_text(shape, "🚫 家长避坑三不要", size=32, bold=True)
    elif shape.placeholder_format.idx == 1:
        tf = shape.text_frame; tf.clear()
        add_para(tf, "不盲目刷题、不强迫书写", size=22, bold=True, space_before=14)
        add_para(tf, "    过早机械训练会扼杀学习兴趣", size=18, space_before=2)
        add_para(tf, "不拿别人家的孩子比较", size=22, bold=True, space_before=14)
        add_para(tf, "    每个孩子都有自己的成长节奏", size=18, space_before=2)
        add_para(tf, "不在孩子面前抱怨学校和老师", size=22, bold=True, space_before=14)
        add_para(tf, "    负面评价会让孩子对小学产生恐惧", size=18, space_before=2)

# ═══════════════════════════════════════════
# Slide 14: 三必须
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[LAYOUTS['标题和内容']])
for shape in slide.placeholders:
    if shape.placeholder_format.idx == 0:
        set_text(shape, "✅ 家长守则三必须", size=32, bold=True)
    elif shape.placeholder_format.idx == 1:
        tf = shape.text_frame; tf.clear()
        add_para(tf, "每天家庭谈话时间", size=22, bold=True, space_before=14)
        add_para(tf, "    倾听孩子在校的见闻和感受", size=18, space_before=2)
        add_para(tf, "鼓励孩子自己解决问题", size=22, bold=True, space_before=14)
        add_para(tf, "    培养独立性，不要包办代替", size=18, space_before=2)
        add_para(tf, "与老师善意、及时沟通", size=22, bold=True, space_before=14)
        add_para(tf, "    家校合作是孩子成长的最佳保障", size=18, space_before=2)

# ═══════════════════════════════════════════
# Slide 15: 章节页 - 走进高照实验学校
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[LAYOUTS['章节页-1_4']])
for shape in slide.placeholders:
    if shape.placeholder_format.idx == 14: set_text(shape, "06", size=48, bold=True)
    elif shape.placeholder_format.idx == 15: set_text(shape, "走进高照实验学校", size=32, bold=True)

# ═══════════════════════════════════════════
# Slide 16: 学校介绍
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[LAYOUTS['标题和内容']])
for shape in slide.placeholders:
    if shape.placeholder_format.idx == 0:
        set_text(shape, "让每一束光都照亮成长", size=32, bold=True)
    elif shape.placeholder_format.idx == 1:
        tf = shape.text_frame; tf.clear()
        add_para(tf, "零起点 · 慢教育 · 长陪伴", size=22, color=RGBColor(0xFF,0x8F,0x00), bold=True, space_before=4)
        add_para(tf, "", size=6, space_before=4)
        add_para(tf, "▎阳光阶梯衔接课程", size=20, bold=True, space_before=8)
        add_para(tf, "    适应周不上文化课", size=18, space_before=2)
        add_para(tf, "▎师资保障", size=20, bold=True, space_before=8)
        add_para(tf, "    双班主任，跟班陪餐", size=18, space_before=2)
        add_para(tf, "▎劳动实践", size=20, bold=True, space_before=8)
        add_para(tf, "    每个孩子都有小岗位", size=18, space_before=2)
        add_para(tf, "▎晚托三档模式", size=20, bold=True, space_before=8)
        add_para(tf, "    基础档 · 拓展档 · 安心档", size=18, space_before=2)
        add_para(tf, "", size=6, space_before=4)
        add_para(tf, "连续三年新生适应期家长满意度 95%+", size=18, color=RGBColor(0x43,0xA0,0x47), bold=True, space_before=8)

# ═══════════════════════════════════════════
# Slide 17: 章节页 - 家校携手
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[LAYOUTS['章节页-1_4']])
for shape in slide.placeholders:
    if shape.placeholder_format.idx == 14: set_text(shape, "07", size=48, bold=True)
    elif shape.placeholder_format.idx == 15: set_text(shape, "家校携手", size=32, bold=True)

# ═══════════════════════════════════════════
# Slide 18: 两栏 - 家校携手
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[LAYOUTS['两栏内容']])
for shape in slide.placeholders:
    ph = shape.placeholder_format.idx
    if ph == 0: set_text(shape, "📋 开学前", size=28, bold=True)
    elif ph == 1:
        tf = shape.text_frame; tf.clear()
        add_para(tf, "• 参加家长会", size=18, space_before=10)
        add_para(tf, "• 填写生活习惯清单", size=18, space_before=10)
        add_para(tf, "• 准备姓名贴", size=18, space_before=10)
    elif ph == 2: set_text(shape, "🏫 开学后", size=28, bold=True)
    elif ph == 3:
        tf = shape.text_frame; tf.clear()
        add_para(tf, "• 按时接送", size=18, space_before=10)
        add_para(tf, "• 关注班级群（静音不刷屏）", size=18, space_before=10)
        add_para(tf, "• 使用家校联系本", size=18, space_before=10)

# ═══════════════════════════════════════════
# Slide 19: 互动答疑
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[LAYOUTS['标题和内容']])
for shape in slide.placeholders:
    if shape.placeholder_format.idx == 0:
        set_text(shape, "互动答疑", size=32, bold=True)
    elif shape.placeholder_format.idx == 1:
        tf = shape.text_frame; tf.clear()
        add_para(tf, "🌱 长期", size=22, bold=True, space_before=4)
        add_para(tf, "• 参加家长开放日", size=18, space_before=8)
        add_para(tf, "• 校长直通车", size=18, space_before=8)
        add_para(tf, "• 积极参与家校活动", size=18, space_before=8)
        add_para(tf, "", size=12, space_before=10)
        add_para(tf, "您的问题，我们认真听。", size=22, color=RGBColor(0x43,0xA0,0x47), bold=True, space_before=10)
        add_para(tf, "🙋  要不要上衔接班？", size=20, space_before=12)
        add_para(tf, "🙋  孩子动作慢怎么办？", size=20, space_before=12)

# ═══════════════════════════════════════════
# Slide 20: 结束页
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[LAYOUTS['结束页-1_4']])
for shape in slide.placeholders:
    if shape.placeholder_format.idx == 14:
        set_text(shape, "慢下来，等一等孩子的节奏", size=36, bold=True, align=PP_ALIGN.CENTER)
    elif shape.placeholder_format.idx == 15:
        set_text(shape, "相信孩子，也相信自己。", size=24, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════
# 删除模板原始封面页 (第一页)
# ═══════════════════════════════════════════
NS = 'http://schemas.openxmlformats.org/presentationml/2006/main'
NSR = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
pres_xml = prs._element
sldIdLst = pres_xml.find('.//{%s}sldIdLst' % NS)

if sldIdLst is not None and len(sldIdLst) > 1:
    first_sldId = sldIdLst[0]
    rId = first_sldId.get('{%s}id' % NSR)
    sldIdLst.remove(first_sldId)
    if rId:
        prs.part.drop_rel(rId)

prs.save(OUT)
print(f"✅ PPT 已保存: {OUT}")
print(f"总页数: {len(prs.slides)}")
