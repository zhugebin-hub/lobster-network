#!/usr/bin/env python3
"""Generate QPPB PPTX presentation with image placeholders"""
from pptx import Presentation
from pptx.util import Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Cm(33.866)
prs.slide_height = Cm(19.05)

DB = RGBColor(0, 51, 102)
MB = RGBColor(0, 102, 153)
LB = RGBColor(200, 230, 255)
WH = RGBColor(255, 255, 255)
BK = RGBColor(0, 0, 0)
GY = RGBColor(128, 128, 128)
LG = RGBColor(240, 240, 240)
RD = RGBColor(204, 0, 0)
TS = 15

def add_run(paragraph, text, size=14, color=BK, bold=False, name='微软雅黑'):
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = name
    return run

def add_tb(slide, left, top, width, height, text, size=14, color=BK, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    add_run(p, text, size, color, bold)
    return txBox

def add_rect(slide, left, top, width, height, fill, text='', size=12, color=BK, bold=False, align=PP_ALIGN.CENTER, border=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if border:
        shape.line.color.rgb = border
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        add_run(p, text, size, color, bold)
    return shape

def add_rounded(slide, left, top, width, height, fill, text='', size=12, color=BK, bold=False, align=PP_ALIGN.CENTER, border=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if border:
        shape.line.color.rgb = border
        shape.line.width = Pt(2)
    else:
        shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        add_run(p, text, size, color, bold)
    return shape

def add_img_ph(slide, left, top, width, height, label='[ 图片占位区 ]'):
    return add_rect(slide, left, top, width, height, LG, label, 14, GY, False, PP_ALIGN.CENTER, RGBColor(180,180,180))

def add_title_bar(slide, title):
    add_rect(slide, Cm(0), Cm(0), Cm(33.866), Cm(2.2), DB)
    add_tb(slide, Cm(1.5), Cm(0.4), Cm(30), Cm(1.5), title, 24, WH, True, PP_ALIGN.LEFT)
    add_rect(slide, Cm(0), Cm(2.2), Cm(33.866), Cm(0.06), MB)

def add_page_num(slide, n):
    add_tb(slide, Cm(30), Cm(18.2), Cm(3.5), Cm(0.6), f'{n} / {TS}', 10, GY, False, PP_ALIGN.RIGHT)

def add_footer(slide):
    add_tb(slide, Cm(1), Cm(18.2), Cm(10), Cm(0.6), 'QPPB技术详解 | 芦熠檑 | 2026', 9, GY)


# ==================== SLIDE 1: Cover ====================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(s, Cm(0), Cm(0), Cm(33.866), Cm(0.3), DB)
add_tb(s, Cm(3), Cm(5.5), Cm(28), Cm(3), 'QPPB技术详解', 40, DB, True, PP_ALIGN.CENTER)
add_tb(s, Cm(3), Cm(9), Cm(28), C