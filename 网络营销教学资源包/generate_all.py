#!/usr/bin/env python3
"""批量生成《网络营销》9个项目完整教学资源包"""

import os
import zipfile
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

BASE_DIR = "/home/admin/.openclaw/workspace/网络营销教学资源包"

# ===== 9个项目定义 =====
PROJECTS = [
    {
        "num": "一", "name": "认识网络营销",
        "subtitle": "网络营销认知与入门",
        "hotspot": "AI营销、抖音电商、直播带货",
        "hours": "4",
        "objectives": "了解网络营销基本概念、特点和发展趋势；认识网络营销与传统营销的区别；了解网络营销的主要模式和工具",
        "key_points": "网络营销的概念与特点；网络营销与传统营销的区别；网络营销的主要模式",
        "difficult_points": "网络营销的思维转变；网络营销模式的分类理解",
        "activities": ["案例分析：抖音电商的崛起", "小组讨论：你身边的网络营销", "实操体验：注册抖音小店"],
        "homework": ["完成学习任务单", "搜集3个网络营销案例并分析", "思考：网络营销对你的专业有什么帮助？"],
        "ppt_pages": 12
    },
    {
        "num": "二", "name": "网络营销环境分析",
        "subtitle": "洞察网络营销环境",
        "hotspot": "直播电商政策、数据安全法、算法推荐规范",
        "hours": "4",
        "objectives": "掌握PEST分析模型；了解网络营销的宏观环境和微观环境；能够运用SWOT分析工具",
        "key_points": "PEST分析模型；网络营销宏观环境要素；SWOT分析方法",
        "difficult_points": "政策环境对网络营销的影响分析；微观环境与宏观环境的区分",
        "activities": ["案例分析：直播电商政策变化", "小组练习：用PEST分析某品牌", "SWOT分析实战"],
        "homework": ["完成学习任务单", "选择一家企业做PEST分析", "搜集最新网络营销相关政策"],
        "ppt_pages": 12
    },
    {
        "num": "三", "name": "网络消费者分析",
        "subtitle": "读懂网络消费者",
        "hotspot": "Z世代消费、种草经济、情绪价值消费",
        "hours": "4",
        "objectives": "了解网络消费者的特征和行为模式；掌握网络消费者购买决策过程；能够进行消费者画像分析",
        "key_points": "网络消费者特征；购买决策过程；消费者画像方法",
        "difficult_points": "网络消费者行为的心理分析；消费者画像的构建方法",
        "activities": ["案例分析：小红书种草逻辑", "小组讨论：Z世代消费特点", "实操：构建消费者画像"],
        "homework": ["完成学习任务单", "分析一个品牌的消费者画像", "调查身边同学的网购习惯"],
        "ppt_pages": 12
    },
    {
        "num": "四", "name": "网络营销调研",
        "subtitle": "网络市场调研方法",
        "hotspot": "大数据调研、问卷星、爬虫数据采集",
        "hours": "4",
        "objectives": "掌握网络市场调研的基本方法；学会设计网络问卷；了解大数据调研工具的使用",
        "key_points": "网络调研方法；问卷设计原则；调研数据分析",
        "difficult_points": "问卷设计的有效性；调研数据的分析和解读",
        "activities": ["案例分析：某品牌调研案例", "实操：用问卷星设计问卷", "小组：收集并分析数据"],
        "homework": ["完成学习任务单", "设计一份网络调研问卷", "收集20份有效问卷并分析"],
        "ppt_pages": 12
    },
    {
        "num": "五", "name": "网络营销策略设计",
        "subtitle": "4P到4C的策略升级",
        "hotspot": "国潮品牌、私域流量、DTC模式",
        "hours": "6",
        "objectives": "掌握网络营销4P策略；理解4C理论；能够设计简单的网络营销策略方案",
        "key_points": "网络营销4P策略；4C理论；策略方案设计方法",
        "difficult_points": "4P到4C的思维转变；策略方案的综合设计",
        "activities": ["案例分析：国潮品牌营销策略", "小组练习：设计4P策略", "方案展示与互评"],
        "homework": ["完成学习任务单", "为一个产品设计网络营销策略", "分析一个国潮品牌的营销策略"],
        "ppt_pages": 14
    },
    {
        "num": "六", "name": "网络广告",
        "subtitle": "精准触达目标用户",
        "hotspot": "信息流广告、KOL投放、程序化购买",
        "hours": "6",
        "objectives": "了解网络广告的主要形式；掌握网络广告的投放策略；学会评估广告效果",
        "key_points": "网络广告形式；广告投放策略；广告效果评估",
        "difficult_points": "信息流广告的投放逻辑；ROI的计算和优化",
        "activities": ["案例分析：抖音信息流广告", "小组讨论：KOL选择标准", "实操：设计广告创意"],
        "homework": ["完成学习任务单", "分析一个网络广告案例", "设计一个信息流广告创意"],
        "ppt_pages": 14
    },
    {
        "num": "七", "name": "搜索引擎营销",
        "subtitle": "SEO与SEM实战",
        "hotspot": "AI搜索优化、百度爱采购、短视频SEO",
        "hours": "6",
        "objectives": "了解搜索引擎营销的基本原理；掌握SEO优化技巧；学会SEM竞价推广",
        "key_points": "SEO优化原理；关键词选择；SEM竞价策略",
        "difficult_points": "SEO长期优化策略；竞价排名的成本控制",
        "activities": ["案例分析：百度SEO实战", "实操：关键词研究", "小组：设计SEM方案"],
        "homework": ["完成学习任务单", "为一个网站做SEO诊断", "设计一个SEM投放方案"],
        "ppt_pages": 14
    },
    {
        "num": "八", "name": "社交媒体营销",
        "subtitle": "玩转社交平台",
        "hotspot": "小红书种草、视频号运营、社群营销",
        "hours": "6",
        "objectives": "了解主流社交媒体平台特点；掌握社交媒体营销方法；学会内容运营和粉丝运营",
        "key_points": "社交媒体平台特点；内容运营方法；粉丝运营策略",
        "difficult_points": "不同平台的运营策略差异；内容创意的持续产出",
        "activities": ["案例分析：小红书爆款笔记", "实操：撰写种草笔记", "小组：制定社群运营方案"],
        "homework": ["完成学习任务单", "撰写一篇小红书种草笔记", "制定一个社群运营方案"],
        "ppt_pages": 14
    },
    {
        "num": "九", "name": "网络营销综合实战",
        "subtitle": "完整营销策划案",
        "hotspot": "AI全链路营销、品牌数字化升级、全域营销",
        "hours": "8",
        "objectives": "综合运用所学知识完成营销策划案；掌握营销策划的基本流程；能够进行方案展示和答辩",
        "key_points": "营销策划流程；方案撰写规范；方案展示技巧",
        "difficult_points": "多策略的综合运用；方案的可执行性评估",
        "activities": ["案例分析：完整营销策划案", "小组：制定策划方案", "方案答辩与点评"],
        "homework": ["完成一份完整的网络营销策划案", "进行方案展示和答辩", "撰写学习反思总结"],
        "ppt_pages": 14
    }
]

# ===== 配色方案 =====
ORANGE = RGBColor(0xFF, 0x6F, 0x00)
BLUE = RGBColor(0x19, 0x76, 0xD2)
GREEN = RGBColor(0x38, 0x8E, 0x3C)
RED = RGBColor(0xD3, 0x2F, 0x2F)
PURPLE = RGBColor(0x7B, 0x1F, 0xA2)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x21, 0x21, 0x21)
GRAY = RGBColor(0x75, 0x75, 0x75)
LIGHT_BG = RGBColor(0xF5, 0xF5, 0xF5)

def add_shape(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_textbox(slide, left, top, width, height, text, size=18, color=DARK, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = "微软雅黑"
    return txBox

def add_bullet_list(slide, left, top, width, height, items, size=16, color=DARK, spacing=Pt(6)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = spacing
        run = p.add_run()
        run.text = f"• {item}"
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.name = "微软雅黑"
    return txBox

def generate_ppt(project):
    """生成单个项目的PPT"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    colors = [ORANGE, BLUE, GREEN, RED, PURPLE]
    color = colors[hash(project["name"]) % len(colors)]
    
    # 第1页：封面
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = WHITE
    add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(3.5), color)
    add_shape(slide, Inches(0), Inches(3.5), Inches(13.333), Inches(0.08), RGBColor(0xFF, 0xD5, 0x4F))
    add_textbox(slide, Inches(3), Inches(0.5), Inches(7.3), Inches(1.5), "📚", size=60, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1.5), Inches(1.5), Inches(10.3), Inches(1.2),
                f"项目{project['num']}：{project['name']}", size=44, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(2), Inches(2.8), Inches(9.3), Inches(0.6),
                project["subtitle"], size=24, color=RGBColor(0xFF, 0xE0, 0xB2), align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(2), Inches(4.5), Inches(9.3), Inches(0.5),
                f"《网络营销》· 高等教育出版社 · 于丽娟 主编", size=18, color=GRAY, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(2), Inches(5.2), Inches(9.3), Inches(0.5),
                f"热点结合：{project['hotspot']}", size=16, color=GRAY, align=PP_ALIGN.CENTER)
    add_shape(slide, Inches(0), Inches(7.0), Inches(13.333), Inches(0.08), RGBColor(0xFF, 0xD5, 0x4F))
    
    # 第2页：学习目标
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background; fill = bg.fill; fill.solid(); fill.fore_color.rgb = WHITE
    add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), color)
    add_textbox(slide, Inches(0.5), Inches(0.25), Inches(12), Inches(0.8), "🎯 学习目标", size=32, color=WHITE, bold=True)
    add_textbox(slide, Inches(0.8), Inches(1.6), Inches(11.7), Inches(0.6),
                f"学时：{project['hours']}学时", size=18, color=GRAY, align=PP_ALIGN.CENTER)
    
    objectives = project["objectives"].split("；")
    for i, obj in enumerate(objectives):
        y = Inches(2.5) + Inches(1.0) * i
        add_shape(slide, Inches(1.0), y, Inches(11.3), Inches(0.8), LIGHT_BG)
        add_textbox(slide, Inches(1.3), y + Inches(0.1), Inches(0.5), Inches(0.6), f"{i+1}", size=24, color=color, bold=True, align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(1.9), y + Inches(0.1), Inches(10), Inches(0.6), obj.strip(), size=18, color=DARK)
    
    # 第3页：教学重难点
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background; fill = bg.fill; fill.solid(); fill.fore_color.rgb = WHITE
    add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), color)
    add_textbox(slide, Inches(0.5), Inches(0.25), Inches(12), Inches(0.8), "📌 教学重难点", size=32, color=WHITE, bold=True)
    
    add_shape(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(5.2), RGBColor(0xE8, 0xF5, 0xE9))
    add_textbox(slide, Inches(1.0), Inches(1.8), Inches(5.1), Inches(0.5), "✅ 教学重点", size=24, color=GREEN, bold=True)
    add_bullet_list(slide, Inches(1.2), Inches(2.4), Inches(4.8), Inches(4.0), project["key_points"].split("；"), size=16, color=DARK)
    
    add_shape(slide, Inches(6.8), Inches(1.6), Inches(5.5), Inches(5.2), RGBColor(0xFF, 0xF3, 0xE0))
    add_textbox(slide, Inches(7.0), Inches(1.8), Inches(5.1), Inches(0.5), "⚠️ 教学难点", size=24, color=ORANGE, bold=True)
    add_bullet_list(slide, Inches(7.2), Inches(2.4), Inches(4.8), Inches(4.0), project["difficult_points"].split("；"), size=16, color=DARK)
    
    # 第4页：导入案例
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background; fill = bg.fill; fill.solid(); fill.fore_color.rgb = WHITE
    add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), color)
    add_textbox(slide, Inches(0.5), Inches(0.25), Inches(12), Inches(0.8), "🔥 热点导入", size=32, color=WHITE, bold=True)
    add_textbox(slide, Inches(0.8), Inches(1.6), Inches(11.7), Inches(0.6),
                f"结合时下热点：{project['hotspot']}", size=20, color=GRAY, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1.0), Inches(2.5), Inches(11.3), Inches(4.0),
                f"【情境创设】\n\n同学们，你们有没有注意到，现在打开手机，到处都是各种推荐和广告？\n\n"
                f"比如你在抖音上刷到一个视频，看完后系统就给你推荐了更多同类内容；\n"
                f"你在淘宝上搜索了一个商品，接下来几天都会看到相关广告……\n\n"
                f"这些现象背后，就是{project['name']}在发挥作用！\n\n"
                f"今天，我们就来学习{project['subtitle']}。",
                size=18, color=DARK)
    
    # 第5-N页：知识讲解（简化版）
    for i in range(5, project["ppt_pages"] - 1):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        bg = slide.background; fill = bg.fill; fill.solid(); fill.fore_color.rgb = WHITE
        add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), color)
        add_textbox(slide, Inches(0.5), Inches(0.25), Inches(12), Inches(0.8), f"📖 知识讲解 {i-4}", size=32, color=WHITE, bold=True)
        add_textbox(slide, Inches(1.0), Inches(1.6), Inches(11.3), Inches(5.0),
                    f"【核心知识点 {i-4}】\n\n"
                    f"（此处为知识讲解页，教师可根据教材内容补充详细知识点）\n\n"
                    f"建议内容：\n"
                    f"• 理论概念讲解\n"
                    f"• 案例分析\n"
                    f"• 互动讨论\n"
                    f"• 实操演示",
                    size=18, color=DARK)
    
    # 倒数第2页：课堂活动
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background; fill = bg.fill; fill.solid(); fill.fore_color.rgb = WHITE
    add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), color)
    add_textbox(slide, Inches(0.5), Inches(0.25), Inches(12), Inches(0.8), "🎮 课堂活动", size=32, color=WHITE, bold=True)
    for i, activity in enumerate(project["activities"]):
        y = Inches(1.6) + Inches(1.5) * i
        add_shape(slide, Inches(1.0), y, Inches(11.3), Inches(1.2), LIGHT_BG)
        add_textbox(slide, Inches(1.3), y + Inches(0.1), Inches(0.5), Inches(0.5), f"{i+1}", size=24, color=color, bold=True, align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(1.9), y + Inches(0.1), Inches(10), Inches(0.5), activity, size=20, color=DARK, bold=True)
        add_textbox(slide, Inches(1.9), y + Inches(0.6), Inches(10), Inches(0.5), "（小组合作，限时讨论，代表发言）", size=14, color=GRAY)
    
    # 最后一页：作业布置
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background; fill = bg.fill; fill.solid(); fill.fore_color.rgb = WHITE
    add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), color)
    add_textbox(slide, Inches(0.5), Inches(0.25), Inches(12), Inches(0.8), "📝 作业布置", size=32, color=WHITE, bold=True)
    for i, hw in enumerate(project["homework"]):
        y = Inches(1.6) + Inches(1.2) * i
        add_shape(slide, Inches(1.0), y, Inches(11.3), Inches(0.9), LIGHT_BG)
        add_textbox(slide, Inches(1.3), y + Inches(0.1), Inches(0.5), Inches(0.5), f"{i+1}", size=24, color=color, bold=True, align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(1.9), y + Inches(0.1), Inches(10), Inches(0.7), hw, size=18, color=DARK)
    
    output_path = os.path.join(BASE_DIR, project["name"], f"{project['name']}_教学PPT.pptx")
    prs.save(output_path)
    print(f"✅ PPT已生成：{output_path}")

def generate_教学设计(project):
    """生成教学设计"""
    content = f"""# 项目{project['num']}：{project['name']} 教学设计

**课程名称：** 《网络营销》  
**项目名称：** 项目{project['num']}：{project['name']}  
**授课学时：** {project['hours']}学时  
**教材：** 高等教育出版社 于丽娟 主编  

---

## 一、教学内容分析

本项目围绕{project['subtitle']}展开，结合时下热点"{project['hotspot']}"，帮助学生掌握网络营销的核心概念和实践方法。

### 教学重点
{project['key_points']}

### 教学难点
{project['difficult_points']}

---

## 二、学情分析

**授课对象：** 中职/高职电商类/商贸类专业学生

**已有基础：**
- 学生日常使用智能手机和社交媒体，对网络营销有感性认识
- 具备一定的互联网使用经验

**学习特点：**
- 对新鲜事物好奇，喜欢动手实践
- 抽象思维能力有待加强，需要借助具体案例
- 部分学生理论基础薄弱，需注重实践操作

---

## 三、教学目标

### 1. 知识目标
{chr(10).join(f"- {obj.strip()}" for obj in project['objectives'].split('；'))}

### 2. 能力目标
- 能够运用所学知识分析网络营销案例
- 能够设计简单的网络营销方案
- 能够使用常见的网络营销工具

### 3. 素养目标
- 培养网络营销思维
- 增强创新意识和实践能力
- 树立诚信营销的理念

---

## 四、教学方法

- **案例分析法：** 结合{project['hotspot']}等时下热点案例
- **任务驱动法：** 通过学习任务单引导学生自主探究
- **小组合作学习：** 分组讨论、协作完成任务
- **体验式学习：** 现场体验网络营销工具

---

## 五、教学过程设计

### 第一课时

| 环节 | 时间 | 教师活动 | 学生活动 | 设计意图 |
|------|------|----------|----------|----------|
| 导入 | 8分钟 | 结合{project['hotspot']}创设情境 | 观看案例，思考回答 | 激发兴趣，引出课题 |
| 新授 | 20分钟 | 讲解核心概念和理论 | 听讲、笔记、提问 | 建立知识框架 |
| 互动 | 10分钟 | 组织小组讨论 | 分组讨论，代表发言 | 加深理解 |
| 小结 | 7分钟 | 总结要点，布置作业 | 记录要点 | 巩固知识 |

### 第二课时（如适用）

| 环节 | 时间 | 教师活动 | 学生活动 | 设计意图 |
|------|------|----------|----------|----------|
| 复习 | 5分钟 | 提问回顾 | 回答 | 温故知新 |
| 深化 | 20分钟 | 深入讲解+案例分析 | 听讲、分析案例 | 提升理解 |
| 实践 | 15分钟 | 指导实操活动 | 动手实践 | 学以致用 |
| 展示 | 3分钟 | 点评总结 | 展示成果 | 展示交流 |

---

## 六、板书设计

```
            项目{project['num']}：{project['name']}

一、核心概念
   {project['key_points'].split('；')[0] if project['key_points'] else '待补充'}

二、关键要点
   {project['key_points']}

三、实践应用
   {project['hotspot']}

四、思考题
   如何运用所学知识解决实际问题？
```

---

## 七、教学评价

| 评价维度 | 评价方式 | 权重 |
|----------|----------|------|
| 知识掌握 | 课堂提问、学习任务单 | 30% |
| 技能实践 | 课堂活动参与度 | 40% |
| 态度表现 | 讨论发言、团队合作 | 30% |

---

## 八、教学反思（课后填写）

_（授课后在此记录教学效果、学生反馈、改进方向等）_

---

**编写人：** _________  
**授课日期：** _________  
**授课班级：** _________
"""
    path = os.path.join(BASE_DIR, project["name"], f"{project['name']}_教学设计.md")
    with open(path, 'w', encoding='utf-8') as f:
        f.write(content)
    print(f"✅ 教学设计已生成：{path}")

def generate_学习任务单(project):
    """生成学习任务单"""
    content = f"""# 项目{project['num']}：{project['name']} 学习任务单

**班级：** ______________  **姓名：** ______________  **日期：** ______________

---

## 🎯 学习目标

1. 了解{project['name']}的基本概念
2. 掌握{project['key_points'].split('；')[0] if project['key_points'] else '核心知识点'}
3. 能够运用所学知识分析实际案例
4. 结合{project['hotspot']}等热点进行思考

---

## 任务一：知识梳理——填一填

### 1. 填空题

（1）网络营销是指 _________________________________________________

（2）网络营销与传统营销的主要区别在于：___________________________

（3）{project['key_points'].split('；')[0] if project['key_points'] else '本项目的核心概念是'}

### 2. 简答题

（1）请简述{project['name']}的核心要点：

_________________________________________________________________

_________________________________________________________________

（2）结合{project['hotspot']}，谈谈你的理解：

_________________________________________________________________

_________________________________________________________________

---

## 任务二：案例分析——析一析

**案例：** 结合{project['hotspot']}选择一个典型案例进行分析

**案例名称：** _________________________________________________

### 分析要点

| 分析维度 | 你的分析 |
|----------|----------|
| 案例背景 | |
| 营销策略 | |
| 成功因素 | |
| 可借鉴之处 | |

**你的思考：**

_________________________________________________________________

_________________________________________________________________

---

## 任务三：课堂活动——做一做

### 活动：{project['activities'][0] if project['activities'] else '小组讨论'}

**小组名称：** ______________  **小组成员：** ______________

**活动任务：** _________________________________________________

**讨论要点：**

1. _____________________________________________________________

2. _____________________________________________________________

3. _____________________________________________________________

**小组结论：**

_________________________________________________________________

_________________________________________________________________

---

## 任务四：拓展思考——想一想

**话题：** 结合{project['hotspot']}，思考{project['name']}在实际中的应用

**你的观点：**

_________________________________________________________________

_________________________________________________________________

_________________________________________________________________

---

## 📋 学习自评

| 评价项目 | 很好 😊 | 一般 😐 | 还需努力 💪 |
|----------|---------|---------|------------|
| 我了解了基本概念 | □ | □ | □ |
| 我掌握了核心要点 | □ | □ | □ |
| 我能分析实际案例 | □ | □ | □ |
| 我参与了课堂活动 | □ | □ | □ |
| 我完成了所有任务 | □ | □ | □ |

**本节课我最大的收获是：**

_________________________________________________________________

**我还想进一步了解的知识：**

_________________________________________________________________

---

**教师评语：**

_________________________________________________________________

**成绩评定：** □ 优秀  □ 良好  □ 合格  □ 待努力
"""
    path = os.path.join(BASE_DIR, project["name"], f"{project['name']}_学习任务单.md")
    with open(path, 'w', encoding='utf-8') as f:
        f.write(content)
    print(f"✅ 学习任务单已生成：{path}")

def generate_作业清单(project):
    """生成作业清单"""
    content = f"""# 项目{project['num']}：{project['name']} 作业清单

**班级：** ______________  **姓名：** ______________  **日期：** ______________

---

## 📝 作业要求

| 作业类型 | 完成时间 | 分值 | 提交方式 |
|----------|----------|------|----------|
| 基础作业 | 课后2天内 | 40分 | 书面提交 |
| 实践作业 | 课后3天内 | 40分 | 电子提交 |
| 拓展作业 | 课后5天内 | 20分 | 口头/书面 |

---

## 一、基础作业（40分）

### 1. 概念理解题（10分）

（1）请用自己的话解释{project['name']}的概念：

_________________________________________________________________

_________________________________________________________________

（2）{project['key_points'].split('；')[0] if project['key_points'] else '请简述本项目的核心概念'}：

_________________________________________________________________

_________________________________________________________________

### 2. 知识应用题（15分）

结合{project['hotspot']}，分析一个网络营销案例：

**案例名称：** _________________________________________________

**案例分析：**

| 分析维度 | 内容 |
|----------|------|
| 案例背景 | |
| 营销策略 | |
| 效果评估 | |
| 你的评价 | |

### 3. 简答题（15分）

{project['difficult_points'].split('；')[0] if project['difficult_points'] else '请简述本项目的难点内容'}：

_________________________________________________________________

_________________________________________________________________

_________________________________________________________________

---

## 二、实践作业（40分）

### 任务：{project['homework'][1] if len(project['homework']) > 1 else '网络营销实践任务'}

**任务要求：**

_________________________________________________________________

_________________________________________________________________

**提交格式：**
- 字数：不少于800字
- 格式：Word文档或PPT
- 内容：包含分析、方案、反思

**评分标准：**

| 评分项 | 分值 | 评分要点 |
|--------|------|----------|
| 内容完整性 | 15分 | 是否覆盖所有要求 |
| 分析深度 | 15分 | 是否有深度思考 |
| 创新性 | 10分 | 是否有创新观点 |

---

## 三、拓展作业（20分）

### 任务：{project['homework'][2] if len(project['homework']) > 2 else '网络营销拓展任务'}

**任务要求：**

_________________________________________________________________

_________________________________________________________________

**可选形式（任选其一）：**
- □ 撰写一篇分析报告
- □ 制作一个PPT展示
- □ 录制一个短视频
- □ 设计一个营销方案

---

## 四、作业提交清单

| 序号 | 作业内容 | 完成状态 | 得分 |
|------|----------|----------|------|
| 1 | 概念理解题 | □ 已完成 □ 未完成 | |
| 2 | 知识应用题 | □ 已完成 □ 未完成 | |
| 3 | 简答题 | □ 已完成 □ 未完成 | |
| 4 | 实践作业 | □ 已完成 □ 未完成 | |
| 5 | 拓展作业 | □ 已完成 □ 未完成 | |

**总分：** ______ / 100分

---

## 五、教师评语

_________________________________________________________________

_________________________________________________________________

**教师签名：** _________  **日期：** _________

---

**温馨提示：**
1. 请按时完成作业，逾期将扣分
2. 实践作业需提交电子文档
3. 如有疑问，请及时向老师请教
4. 鼓励创新思考，不抄袭、不雷同
"""
    path = os.path.join(BASE_DIR, project["name"], f"{project['name']}_作业清单.md")
    with open(path, 'w', encoding='utf-8') as f:
        f.write(content)
    print(f"✅ 作业清单已生成：{path}")

# ===== 主程序 =====
if __name__ == "__main__":
    print("=" * 60)
    print("🚀 开始批量生成《网络营销》9个项目教学资源包")
    print("=" * 60)
    
    for project in PROJECTS:
        print(f"\n📦 正在生成项目{project['num']}：{project['name']}...")
        
        # 创建项目目录
        project_dir = os.path.join(BASE_DIR, project["name"])
        os.makedirs(project_dir, exist_ok=True)
        
        # 生成PPT
        generate_ppt(project)
        
        # 生成教学设计
        generate_教学设计(project)
        
        # 生成学习任务单
        generate_学习任务单(project)
        
        # 生成作业清单
        generate_作业清单(project)
        
        print(f"✅ 项目{project['num']}：{project['name']} 完成！")
    
    print("\n" + "=" * 60)
    print("🎉 所有项目资源包生成完成！")
    print("=" * 60)
