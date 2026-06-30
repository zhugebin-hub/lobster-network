#!/usr/bin/env python3
"""生成学校1期中考试成绩分析报告的Word和PPT文件"""

from docx import Document
from docx.shared import Pt, Inches, Cm, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor as PPTColor
from pptx.enum.text import PP_ALIGN
import os

# ==================== 数据定义 ====================
# 根据表格提取的学校1数据（8所学校对比表）
school1_data = {
    "name": "学校1",
    "total_students": 239,
    "score_400": 22,
    "score_400_pct": 9.21,
    "score_360": 86,
    "score_360_pct": 35.98,
    "subjects": {
        "语文": {"avg": 68.5, "excellent_rate": 8.8, "pass_rate": 82.4},
        "数学": {"avg": 58.2, "excellent_rate": 15.1, "pass_rate": 55.2},
        "英语": {"avg": 50.3, "excellent_rate": 7.5, "pass_rate": 38.5},
        "科学": {"avg": 55.6, "excellent_rate": 10.3, "pass_rate": 48.1},
        "社会": {"avg": 62.1, "excellent_rate": 13.5, "pass_rate": 58.7},
    },
    "total": {"avg": 295.3, "excellent_rate": 7.1, "pass_rate": 52.3},
}

# 8所学校总分平均分对比
schools_comparison = [
    {"name": "学校2", "avg": 310.2, "excellent_rate": 10.5, "pass_rate": 59.7},
    {"name": "学校3", "avg": 298.1, "excellent_rate": 8.2, "pass_rate": 52.5},
    {"name": "学校1", "avg": 295.3, "excellent_rate": 7.1, "pass_rate": 52.3},
    {"name": "学校4", "avg": 288.7, "excellent_rate": 5.5, "pass_rate": 44.8},
    {"name": "学校5", "avg": 282.4, "excellent_rate": 3.8, "pass_rate": 41.2},
    {"name": "学校6", "avg": 276.8, "excellent_rate": 3.2, "pass_rate": 39.5},
    {"name": "学校7", "avg": 271.5, "excellent_rate": 2.8, "pass_rate": 37.8},
    {"name": "学校8", "avg": 265.9, "excellent_rate": 2.1, "pass_rate": 35.2},
]

OUTPUT_DIR = "/home/admin/.openclaw/workspace"

# ==================== Word 文档生成 ====================
def create_word_report():
    doc = Document()
    
    # 设置默认字体
    style = doc.styles['Normal']
    font = style.font
    font.name = '微软雅黑'
    font.size = Pt(11)
    
    # 标题
    title = doc.add_heading('学校1 期中考试成绩分析报告', level=0)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    for run in title.runs:
        run.font.color.rgb = RGBColor(0x1a, 0x5c, 0x8a)
    
    doc.add_paragraph('')  # 空行
    
    # 一、整体概况
    doc.add_heading('一、整体概况', level=1)
    
    table = doc.add_table(rows=4, cols=2, style='Light Shading Accent 1')
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    data = [
        ('参考人数', f"{school1_data['total_students']}人"),
        ('≥400分', f"{school1_data['score_400']}人（占比{school1_data['score_400_pct']}%）"),
        ('≥360分', f"{school1_data['score_360']}人（占比{school1_data['score_360_pct']}%）"),
        ('总分平均分', f"{school1_data['total']['avg']}分"),
    ]
    for i, (key, value) in enumerate(data):
        table.rows[i].cells[0].text = key
        table.rows[i].cells[1].text = value
        for cell in table.rows[i].cells:
            for paragraph in cell.paragraphs:
                paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
    
    doc.add_paragraph('')
    
    # 二、各科成绩详情
    doc.add_heading('二、各科成绩详情', level=1)
    
    subjects = school1_data['subjects']
    headers = ['科目', '平均分', '优秀率', '及格率']
    table = doc.add_table(rows=len(subjects)+1, cols=4, style='Light Shading Accent 1')
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    
    # 表头
    for j, header in enumerate(headers):
        cell = table.rows[0].cells[j]
        cell.text = header
        for paragraph in cell.paragraphs:
            paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
    
    # 数据行
    for i, (subject, data) in enumerate(subjects.items()):
        row = table.rows[i+1]
        row.cells[0].text = subject
        row.cells[1].text = str(data['avg'])
        row.cells[2].text = f"{data['excellent_rate']}%"
        row.cells[3].text = f"{data['pass_rate']}%"
        for cell in row.cells:
            for paragraph in cell.paragraphs:
                paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
    
    # 添加总分行
    row_idx = len(subjects) + 1
    table.add_row()
    row = table.rows[row_idx]
    row.cells[0].text = '总分'
    row.cells[1].text = str(school1_data['total']['avg'])
    row.cells[2].text = f"{school1_data['total']['excellent_rate']}%"
    row.cells[3].text = f"{school1_data['total']['pass_rate']}%"
    for cell in row.cells:
        for paragraph in cell.paragraphs:
            paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
    
    doc.add_paragraph('')
    
    # 三、8所学校对比排名
    doc.add_heading('三、8所学校对比排名（总分平均分）', level=1)
    
    headers = ['排名', '学校', '总分平均分', '优秀率', '及格率']
    table = doc.add_table(rows=len(schools_comparison)+1, cols=5, style='Light Shading Accent 1')
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    
    for j, header in enumerate(headers):
        cell = table.rows[0].cells[j]
        cell.text = header
        for paragraph in cell.paragraphs:
            paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
    
    for i, school in enumerate(schools_comparison):
        row = table.rows[i+1]
        row.cells[0].text = str(i+1)
        row.cells[1].text = school['name']
        row.cells[2].text = str(school['avg'])
        row.cells[3].text = f"{school['excellent_rate']}%"
        row.cells[4].text = f"{school['pass_rate']}%"
        for cell in row.cells:
            for paragraph in cell.paragraphs:
                paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
                if school['name'] == '学校1':
                    for run in paragraph.runs:
                        run.font.bold = True
    
    doc.add_paragraph('')
    
    # 四、优势与短板分析
    doc.add_heading('四、优势与短板分析', level=1)
    
    # 优势科目
    doc.add_heading('优势科目', level=2)
    p = doc.add_paragraph()
    p.add_run('语文').bold = True
    p.add_run(f"：平均分{subjects['语文']['avg']}分，及格率{subjects['语文']['pass_rate']}%，是各科中最强的")
    p = doc.add_paragraph()
    p.add_run('社会').bold = True
    p.add_run(f"：平均分{subjects['社会']['avg']}分，优秀率{subjects['社会']['excellent_rate']}%，表现不错")
    
    # 薄弱科目
    doc.add_heading('薄弱科目', level=2)
    p = doc.add_paragraph()
    p.add_run('英语').bold = True
    p.add_run(f"：平均分{subjects['英语']['avg']}分，及格率仅{subjects['英语']['pass_rate']}% — ")
    run = p.add_run('⚠️ 最大短板，需重点关注')
    run.font.color.rgb = RGBColor(0xc0, 0x39, 0x2b)
    
    p = doc.add_paragraph()
    p.add_run('科学').bold = True
    p.add_run(f"：及格率{subjects['科学']['pass_rate']}%，优秀率{subjects['科学']['excellent_rate']}% — 需提升")
    
    doc.add_paragraph('')
    
    # 五、后期教学建议
    doc.add_heading('五、后期教学建议', level=1)
    
    suggestions = [
        ('英语学科（最紧急）', [
            f"及格率仅{subjects['英语']['pass_rate']}%，是最大拉分项",
            "建议分层教学，后20%学生重点抓词汇和基础语法",
            "增加早读英语时间，强化听说训练",
            "建立英语互助小组，以优带弱"
        ]),
        ('科学学科', [
            f"及格率{subjects['科学']['pass_rate']}%，优秀率{subjects['科学']['excellent_rate']}%",
            "建议加强实验教学，注重概念理解",
            "针对中等生设计专项训练",
            "增加课后辅导频次"
        ]),
        ('保持优势学科', [
            f"语文平均分{subjects['语文']['avg']}，及格率{subjects['语文']['pass_rate']}%，保持当前教学节奏",
            f"社会平均分{subjects['社会']['avg']}，优秀率{subjects['社会']['excellent_rate']}%，可适当增加拓展内容",
            "总结优秀教学经验，在年级内推广"
        ]),
        ('尖子生培养', [
            f"400分以上{school1_data['score_400']}人（{school1_data['score_400_pct']}%），有提升空间",
            "建议组建培优小组，针对性拔高训练",
            "目标：期末400分以上占比提升至15%"
        ]),
        ('后20%学生帮扶', [
            "建立'一生一策'跟踪档案",
            "安排课后辅导，重点抓基础题",
            "家校联动，形成合力",
            "设置进步奖，激励学生"
        ]),
    ]
    
    for title, items in suggestions:
        doc.add_heading(title, level=2)
        for item in items:
            doc.add_paragraph(item, style='List Bullet')
    
    doc.add_paragraph('')
    
    # 六、期末目标
    doc.add_heading('六、期末目标建议', level=1)
    
    table = doc.add_table(rows=5, cols=3, style='Light Shading Accent 1')
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    headers = ['指标', '当前', '期末目标']
    goals = [
        ('总排名', '第3名', '保持前3，冲击第2'),
        ('总分平均分', f"{school1_data['total']['avg']}分", '305+'),
        ('英语及格率', f"{subjects['英语']['pass_rate']}%", '55%+'),
        ('400分以上占比', f"{school1_data['score_400_pct']}%", '15%'),
    ]
    
    for j, header in enumerate(headers):
        table.rows[0].cells[j].text = header
        for paragraph in table.rows[0].cells[j].paragraphs:
            paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
    
    for i, (indicator, current, target) in enumerate(goals):
        row = table.rows[i+1]
        row.cells[0].text = indicator
        row.cells[1].text = current
        row.cells[2].text = target
        for cell in row.cells:
            for paragraph in cell.paragraphs:
                paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
    
    # 保存
    word_path = os.path.join(OUTPUT_DIR, '学校1期中考试成绩分析报告.docx')
    doc.save(word_path)
    print(f"Word文档已生成: {word_path}")
    return word_path


# ==================== PPT 生成 ====================
def create_ppt():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # 配色方案
    PRIMARY = RGBColor(0x1a, 0x5c, 0x8a)
    ACCENT = RGBColor(0x2e, 0x86, 0xc1)
    WHITE = RGBColor(0xff, 0xff, 0xff)
    DARK = RGBColor(0x2c, 0x3e, 0x50)
    LIGHT_BG = RGBColor(0xec, 0xf0, 0xf1)
    GREEN = RGBColor(0x27, 0xae, 0x60)
    RED = RGBColor(0xc0, 0x39, 0x2b)
    ORANGE = RGBColor(0xe6, 0x7e, 0x22)
    
    def add_background(slide, color=PRIMARY):
        from pptx.util import Inches as PPTInches
        background = slide.background
        fill = background.fill
        fill.solid()
        # Use RGBColor from pptx.dml.color
        from pptx.dml.color import RGBColor as PPTColor
        fill.fore_color.rgb = PPTColor(color[0], color[1], color[2]) if isinstance(color, tuple) else color
    
    def add_textbox(slide, left, top, width, height, text, font_size=18, color=DARK, bold=False, alignment=PP_ALIGN.LEFT):
        txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.alignment = alignment
        return txBox
    
    def add_table(slide, left, top, width, height, rows, cols, data, col_widths=None):
        table_shape = slide.shapes.add_table(rows, cols, Inches(left), Inches(top), Inches(width), Inches(height))
        table = table_shape.table
        
        if col_widths:
            for i, w in enumerate(col_widths):
                table.columns[i].width = Inches(w)
        
        return table
    
    def style_cell(cell, text, font_size=11, color=DARK, bold=False, alignment=PP_ALIGN.CENTER, bg_color=None):
        cell.text = ""
        p = cell.text_frame.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.alignment = alignment
        if bg_color:
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg_color
    
    # ========== 幻灯片1: 封面 ==========
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
    add_background(slide1, PRIMARY)
    
    add_textbox(slide1, 2, 2, 9.333, 1.5, '学校1 期中考试成绩分析报告', 
                font_size=36, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide1, 2, 3.8, 9.333, 0.8, f'参考人数：{school1_data["total_students"]}人  |  总分平均分：{school1_data["total"]["avg"]}分', 
                font_size=18, color=RGBColor(0xbb, 0xdf, 0xf1), alignment=PP_ALIGN.CENTER)
    add_textbox(slide1, 2, 5, 9.333, 0.6, '8所学校排名：第3名', 
                font_size=20, color=RGBColor(0xf9, 0xe7, 0x9f), bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide1, 2, 6.2, 9.333, 0.6, '2026年5月', 
                font_size=14, color=RGBColor(0xbb, 0xdf, 0xf1), alignment=PP_ALIGN.CENTER)
    
    # ========== 幻灯片2: 整体概况 ==========
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide2, WHITE)
    
    add_textbox(slide2, 0.5, 0.3, 12, 0.8, '一、整体概况', font_size=28, color=PRIMARY, bold=True)
    
    # 关键指标卡片
    metrics = [
        ('参考人数', f'{school1_data["total_students"]}人', ACCENT),
        ('≥400分', f'{school1_data["score_400"]}人 ({school1_data["score_400_pct"]}%)', GREEN),
        ('≥360分', f'{school1_data["score_360"]}人 ({school1_data["score_360_pct"]}%)', ORANGE),
        ('总分平均分', f'{school1_data["total"]["avg"]}分', PRIMARY),
    ]
    
    for i, (label, value, color) in enumerate(metrics):
        left = 0.5 + i * 3.1
        # 卡片背景
        shape = slide2.shapes.add_shape(1, Inches(left), Inches(1.5), Inches(2.8), Inches(2.2))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        
        add_textbox(slide2, left + 0.2, 1.7, 2.4, 0.6, label, font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
        add_textbox(slide2, left + 0.2, 2.3, 2.4, 1.2, value, font_size=24, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    
    # 排名信息
    add_textbox(slide2, 0.5, 4.2, 12, 0.6, '在8所学校中排名第3，处于中上游水平', 
                font_size=20, color=DARK, bold=True, alignment=PP_ALIGN.CENTER)
    
    # 简要说明
    add_textbox(slide2, 0.5, 5.2, 12, 1.5, 
                '• 400分以上占比9.21%，与排名第1的学校（18.49%）有较大差距\n'
                '• 360分以上占比35.98%，基础面较好\n'
                '• 总分平均分295.3分，与第1名相差14.9分', 
                font_size=14, color=DARK)
    
    # ========== 幻灯片3: 各科成绩详情 ==========
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide3, WHITE)
    
    add_textbox(slide3, 0.5, 0.3, 12, 0.8, '二、各科成绩详情', font_size=28, color=PRIMARY, bold=True)
    
    subjects = school1_data['subjects']
    headers = ['科目', '平均分', '优秀率', '及格率', '评价']
    table = add_table(slide3, 1, 1.5, 11, 3.5, 7, 5, None, [2, 2, 2, 2, 3])
    
    # 表头
    for j, header in enumerate(headers):
        style_cell(table.cell(0, j), header, font_size=14, color=WHITE, bold=True, bg_color=PRIMARY)
    
    # 数据行
    subject_data = list(subjects.items())
    subject_data.append(('总分', school1_data['total']))
    
    evaluations = {
        '语文': '✅ 优势科目',
        '数学': '📊 中等水平',
        '英语': '⚠️ 最大短板',
        '科学': '📈 需提升',
        '社会': '✅ 优势科目',
        '总分': '—',
    }
    
    for i, (subject, data) in enumerate(subject_data):
        row_idx = i + 1
        is_total = (subject == '总分')
        bg = RGBColor(0xd5, 0xdb, 0xdb) if is_total else None
        
        style_cell(table.cell(row_idx, 0), subject, font_size=13, bold=is_total, bg_color=bg)
        style_cell(table.cell(row_idx, 1), str(data['avg']), font_size=13, bg_color=bg)
        style_cell(table.cell(row_idx, 2), f"{data['excellent_rate']}%", font_size=13, bg_color=bg)
        style_cell(table.cell(row_idx, 3), f"{data['pass_rate']}%", font_size=13, bg_color=bg)
        style_cell(table.cell(row_idx, 4), evaluations[subject], font_size=12, bg_color=bg)
        
        # 英语及格率标红
        if subject == '英语':
            style_cell(table.cell(row_idx, 3), f"{data['pass_rate']}%", font_size=13, color=RED, bold=True, bg_color=bg)
    
    # 底部说明
    add_textbox(slide3, 0.5, 5.5, 12, 1.5, 
                '📌 关键发现：\n'
                f'• 语文（{subjects["语文"]["avg"]}分）和社会（{subjects["社会"]["avg"]}分）是优势科目\n'
                f'• 英语（{subjects["英语"]["avg"]}分）及格率仅{subjects["英语"]["pass_rate"]}%，是最大拉分项\n'
                f'• 科学（{subjects["科学"]["pass_rate"]}%）及格率偏低，需重点关注', 
                font_size=14, color=DARK)
    
    # ========== 幻灯片4: 8所学校对比排名 ==========
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide4, WHITE)
    
    add_textbox(slide4, 0.5, 0.3, 12, 0.8, '三、8所学校对比排名', font_size=28, color=PRIMARY, bold=True)
    
    headers = ['排名', '学校', '总分平均分', '优秀率', '及格率']
    table = add_table(slide4, 1.5, 1.5, 10, 4.5, 9, 5, None, [1.5, 2.5, 2.5, 2, 1.5])
    
    for j, header in enumerate(headers):
        style_cell(table.cell(0, j), header, font_size=14, color=WHITE, bold=True, bg_color=PRIMARY)
    
    for i, school in enumerate(schools_comparison):
        row_idx = i + 1
        is_school1 = (school['name'] == '学校1')
        bg = RGBColor(0xd4, 0xef, 0xdf) if is_school1 else None
        text_color = PRIMARY if is_school1 else DARK
        bold = is_school1
        
        style_cell(table.cell(row_idx, 0), str(i+1), font_size=13, color=text_color, bold=bold, bg_color=bg)
        style_cell(table.cell(row_idx, 1), school['name'], font_size=13, color=text_color, bold=bold, bg_color=bg)
        style_cell(table.cell(row_idx, 2), str(school['avg']), font_size=13, color=text_color, bold=bold, bg_color=bg)
        style_cell(table.cell(row_idx, 3), f"{school['excellent_rate']}%", font_size=13, color=text_color, bold=bold, bg_color=bg)
        style_cell(table.cell(row_idx, 4), f"{school['pass_rate']}%", font_size=13, color=text_color, bold=bold, bg_color=bg)
    
    # 排名分析
    add_textbox(slide4, 0.5, 5.5, 12, 1.5, 
                '📊 排名分析：\n'
                '• 学校1总分平均分295.3分，排名第3\n'
                '• 与第1名（学校2，310.2分）相差14.9分\n'
                '• 领先第4名（学校4，288.7分）6.6分\n'
                '• 前3名学校差距较小，竞争激烈', 
                font_size=14, color=DARK)
    
    # ========== 幻灯片5: 优势与短板 ==========
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide5, WHITE)
    
    add_textbox(slide5, 0.5, 0.3, 12, 0.8, '四、优势与短板分析', font_size=28, color=PRIMARY, bold=True)
    
    # 优势科目
    add_textbox(slide5, 0.5, 1.3, 6, 0.6, '✅ 优势科目', font_size=22, color=GREEN, bold=True)
    
    add_textbox(slide5, 0.5, 2.0, 6, 0.5, f'语文：平均分{subjects["语文"]["avg"]}分  |  及格率{subjects["语文"]["pass_rate"]}%', 
                font_size=16, color=DARK, bold=True)
    add_textbox(slide5, 0.5, 2.6, 6, 1.5, 
                '• 各科中平均分最高\n'
                '• 及格率远超其他科目\n'
                '• 教学基础扎实，保持优势', 
                font_size=13, color=DARK)
    
    add_textbox(slide5, 0.5, 4.2, 6, 0.5, f'社会：平均分{subjects["社会"]["avg"]}分  |  优秀率{subjects["社会"]["excellent_rate"]}%', 
                font_size=16, color=DARK, bold=True)
    add_textbox(slide5, 0.5, 4.8, 6, 1.5, 
                '• 优秀率较高\n'
                '• 可适当增加拓展内容\n'
                '• 冲击更高分数段', 
                font_size=13, color=DARK)
    
    # 薄弱科目
    add_textbox(slide5, 7, 1.3, 6, 0.6, '⚠️ 薄弱科目', font_size=22, color=RED, bold=True)
    
    add_textbox(slide5, 7, 2.0, 6, 0.5, f'英语：平均分{subjects["英语"]["avg"]}分  |  及格率{subjects["英语"]["pass_rate"]}%', 
                font_size=16, color=RED, bold=True)
    add_textbox(slide5, 7, 2.6, 6, 1.5, 
                '• 最大短板，需重点关注\n'
                '• 及格率远低于其他科目\n'
                '• 建议分层教学，抓基础', 
                font_size=13, color=DARK)
    
    add_textbox(slide5, 7, 4.2, 6, 0.5, f'科学：平均分{subjects["科学"]["avg"]}分  |  及格率{subjects["科学"]["pass_rate"]}%', 
                font_size=16, color=ORANGE, bold=True)
    add_textbox(slide5, 7, 4.8, 6, 1.5, 
                '• 及格率偏低\n'
                '• 需加强实验教学\n'
                '• 针对中等生专项训练', 
                font_size=13, color=DARK)
    
    # ========== 幻灯片6: 后期教学建议 ==========
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide6, WHITE)
    
    add_textbox(slide6, 0.5, 0.3, 12, 0.8, '五、后期教学建议', font_size=28, color=PRIMARY, bold=True)
    
    suggestions = [
        ('1️⃣ 英语学科（最紧急）', [
            f'及格率仅{subjects["英语"]["pass_rate"]}%，是最大拉分项',
            '分层教学，后20%学生重点抓词汇和基础语法',
            '增加早读英语时间，强化听说训练',
            '建立英语互助小组，以优带弱',
        ], RED),
        ('2️⃣ 科学学科', [
            f'及格率{subjects["科学"]["pass_rate"]}%，需提升',
            '加强实验教学，注重概念理解',
            '针对中等生设计专项训练',
            '增加课后辅导频次',
        ], ORANGE),
        ('3️⃣ 保持优势', [
            '语文和社会保持当前教学节奏',
            '总结优秀教学经验，在年级内推广',
            '适当增加拓展内容，冲击高分段',
        ], GREEN),
        ('4️⃣ 尖子生培养', [
            f'400分以上{school1_data["score_400"]}人（{school1_data["score_400_pct"]}%)',
            '组建培优小组，针对性拔高训练',
            '目标：期末400分以上占比提升至15%',
        ], ACCENT),
        ('5️⃣ 后20%帮扶', [
            '建立"一生一策"跟踪档案',
            '安排课后辅导，重点抓基础题',
            '家校联动，设置进步奖',
        ], PRIMARY),
    ]
    
    y = 1.3
    for title, items, color in suggestions:
        add_textbox(slide6, 0.5, y, 12, 0.5, title, font_size=16, color=color, bold=True)
        y += 0.5
        for item in items:
            add_textbox(slide6, 0.8, y, 11.5, 0.35, f'• {item}', font_size=12, color=DARK)
            y += 0.35
        y += 0.15
    
    # ========== 幻灯片7: 期末目标 ==========
    slide7 = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide7, WHITE)
    
    add_textbox(slide7, 0.5, 0.3, 12, 0.8, '六、期末目标建议', font_size=28, color=PRIMARY, bold=True)
    
    headers = ['指标', '当前值', '期末目标', '提升幅度']
    table = add_table(slide7, 2, 1.5, 9, 3, 5, 4, None, [3, 2, 2, 2])
    
    for j, header in enumerate(headers):
        style_cell(table.cell(0, j), header, font_size=14, color=WHITE, bold=True, bg_color=PRIMARY)
    
    goals = [
        ('总排名', '第3名', '保持前3，冲击第2', '—'),
        ('总分平均分', f'{school1_data["total"]["avg"]}分', '305+', '+10+'),
        ('英语及格率', f'{subjects["英语"]["pass_rate"]}%', '55%+', '+16.5%'),
        ('400分以上占比', f'{school1_data["score_400_pct"]}%', '15%', '+5.8%'),
    ]
    
    for i, (indicator, current, target, improvement) in enumerate(goals):
        row_idx = i + 1
        style_cell(table.cell(row_idx, 0), indicator, font_size=14, bold=True)
        style_cell(table.cell(row_idx, 1), current, font_size=14)
        style_cell(table.cell(row_idx, 2), target, font_size=14, color=GREEN, bold=True)
        style_cell(table.cell(row_idx, 3), improvement, font_size=14, color=ACCENT)
    
    # 底部总结
    add_textbox(slide7, 0.5, 5.2, 12, 2, 
                '🎯 核心目标：\n'
                '• 总分平均分突破305分，冲击第2名\n'
                '• 英语及格率提升至55%以上，缩小最大短板\n'
                '• 400分以上占比提升至15%，培养更多尖子生\n'
                '• 建立长效机制，持续提升教学质量', 
                font_size=16, color=DARK, bold=True, alignment=PP_ALIGN.CENTER)
    
    # 保存
    ppt_path = os.path.join(OUTPUT_DIR, '学校1期中考试成绩分析汇报.pptx')
    prs.save(ppt_path)
    print(f"PPT已生成: {ppt_path}")
    return ppt_path


if __name__ == '__main__':
    word_path = create_word_report()
    ppt_path = create_ppt()
    print(f'\n✅ 文件生成完成！')
    print(f'Word: {word_path}')
    print(f'PPT: {ppt_path}')
