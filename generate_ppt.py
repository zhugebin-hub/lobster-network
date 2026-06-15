#!/usr/bin/env python3
"""Generate PPT from 嘉善信息技术工程学校校园学生突发事件应急预案"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import copy

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color palette
DARK_RED = RGBColor(0x8B, 0x00, 0x00)
RED = RGBColor(0xC4, 0x1E, 0x3A)
LIGHT_RED = RGBColor(0xFF, 0xE4, 0xE1)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
MED_GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
ACCENT_RED = RGBColor(0xB2, 0x22, 0x22)
DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)

def add_background(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape_bg(slide, left, top, width, height, color, alpha=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if alpha is not None:
        from lxml import etree
        nsmap = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
        solidFill = shape.fill._fill
        srgb = solidFill.find('.//a:srgbClr', nsmap)
        if srgb is not None:
            alpha_elem = etree.SubElement(srgb, '{http://schemas.openxmlformats.org/drawingml/2006/main}alpha')
            alpha_elem.set('val', str(int(alpha * 1000)))
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18, bold=False, color=DARK_GRAY, alignment=PP_ALIGN.LEFT, font_name='微软雅黑'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_bullet_list(slide, left, top, width, height, items, font_size=16, color=DARK_GRAY, bullet_color=RED, spacing=Pt(6), font_name='微软雅黑'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = spacing
        p.level = 0
        # Add bullet
        pPr = p._pPr
        if pPr is None:
            from lxml import etree
            nsmap = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
            pPr = etree.SubElement(p._p, '{http://schemas.openxmlformats.org/drawingml/2006/main}pPr')
        from lxml import etree
        nsmap_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
        buChar = etree.SubElement(pPr, f'{{{nsmap_a}}}buChar')
        buChar.set('char', '●')
        buClr = etree.SubElement(pPr, f'{{{nsmap_a}}}buClr')
        srgb = etree.SubElement(buClr, f'{{{nsmap_a}}}srgbClr')
        srgb.set('val', f'{bullet_color[0]:02X}{bullet_color[1]:02X}{bullet_color[2]:02X}')
    return txBox

def add_section_header_bar(slide, text, y_pos=Inches(0.3)):
    # Top red bar
    bar = add_shape_bg(slide, Inches(0), y_pos, Inches(13.333), Inches(0.08), RED)
    # Section title
    add_text_box(slide, Inches(0.8), y_pos + Inches(0.15), Inches(11), Inches(0.6),
                 text, font_size=28, bold=True, color=DARK_RED)

def add_footer(slide, text="嘉善信息技术工程学校 校园学生突发事件应急预案"):
    add_text_box(slide, Inches(0.5), Inches(7.0), Inches(12), Inches(0.4),
                 text, font_size=10, color=MED_GRAY, alignment=PP_ALIGN.CENTER)


# ============================================
# SLIDE 1: Title Slide
# ============================================
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
add_background(slide1, DARK_BG)

# Decorative top bar
add_shape_bg(slide1, Inches(0), Inches(0), Inches(13.333), Inches(0.15), RED)

# Main title
add_text_box(slide1, Inches(1.5), Inches(1.8), Inches(10), Inches(1.2),
             "校园学生突发事件应急预案", font_size=44, bold=True, color=WHITE,
             alignment=PP_ALIGN.CENTER, font_name='微软雅黑')

# Subtitle
add_text_box(slide1, Inches(1.5), Inches(3.2), Inches(10), Inches(0.8),
             "嘉善信息技术工程学校", font_size=28, bold=False, color=RGBColor(0xCC, 0xCC, 0xCC),
             alignment=PP_ALIGN.CENTER, font_name='微软雅黑')

# Date
add_text_box(slide1, Inches(1.5), Inches(4.2), Inches(10), Inches(0.6),
             "（2025年9月修订）", font_size=22, bold=False, color=RGBColor(0x99, 0x99, 0x99),
             alignment=PP_ALIGN.CENTER, font_name='微软雅黑')

# Decorative bottom bar
add_shape_bg(slide1, Inches(0), Inches(7.35), Inches(13.333), Inches(0.15), RED)


# ============================================
# SLIDE 2: Purpose / Overview
# ============================================
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide2, WHITE)
add_section_header_bar(slide2, "一、总体描述 — 编制目的")
add_footer(slide2)

# Content box
add_shape_bg(slide2, Inches(0.8), Inches(1.3), Inches(11.7), Inches(5.3), LIGHT_GRAY)

add_text_box(slide2, Inches(1.2), Inches(1.6), Inches(10.9), Inches(4.5),
             "为确保发生学生突发事件时各项应急工作高效、有序地进行，\n"
             "最大限度地减少人员伤亡和财产损失，\n"
             "稳定社会秩序和校园秩序，\n\n"
             "根据县政府和县教育局有关文件精神，\n"
             "结合我校工作实际，特制定本预案。",
             font_size=22, color=DARK_GRAY, alignment=PP_ALIGN.LEFT, font_name='微软雅黑')


# ============================================
# SLIDE 3: Organization & Leadership
# ============================================
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide3, WHITE)
add_section_header_bar(slide3, "一、总体描述 — 组织与领导")
add_footer(slide3)

add_shape_bg(slide3, Inches(0.8), Inches(1.3), Inches(11.7), Inches(5.3), LIGHT_GRAY)

items = [
    "校长、书记是应急预案的正、副总指挥",
    "应急预案小组成员分别担任各工作组负责人或成员",
    "运行原则：统一指挥、分级响应、岗位责任、互相配合",
    "采取一切必要手段，组织全面救护，将损失降到最低",
    "调动积极因素，做好稳定教育教学秩序和善后安抚工作",
    "对应急工作中发生的争议采取紧急处理措施",
    "向上级部门通报应急救援行动方案，请求支援",
]
add_bullet_list(slide3, Inches(1.2), Inches(1.5), Inches(10.9), Inches(5.0),
                items, font_size=18, color=DARK_GRAY, spacing=Pt(12))


# ============================================
# SLIDE 4: Prevention Guidelines
# ============================================
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide4, WHITE)
add_section_header_bar(slide4, "一、总体描述 — 贯彻预防为主的方针")
add_footer(slide4)

add_shape_bg(slide4, Inches(0.8), Inches(1.3), Inches(11.7), Inches(5.3), LIGHT_GRAY)

items = [
    "加强领导，健全组织，制定应急预案和落实各项措施",
    "校安处做好预案的发放、登记、修改和定期组织学习",
    "加强教师师德规范教育，增强责任意识和法制意识",
    "加强学生法纪教育、安全教育，增强自我保护意识",
    "后勤部门确保校医室常备应急医疗物资，确保医生值班",
    "健全学校各项规章制度，集体性活动要有教师带班",
    "行政值日履行值日工作职责，坚守学校",
    "确保所有行政、班主任、值班人员24小时电话畅通",
]
add_bullet_list(slide4, Inches(1.2), Inches(1.5), Inches(10.9), Inches(5.0),
                items, font_size=17, color=DARK_GRAY, spacing=Pt(10))


# ============================================
# SLIDE 5: Activation Conditions
# ============================================
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide5, WHITE)
add_section_header_bar(slide5, "一、总体描述 — 启用应急预案的情况")
add_footer(slide5)

# Two big cards
add_shape_bg(slide5, Inches(1.0), Inches(2.0), Inches(5.0), Inches(3.5), LIGHT_RED)
add_text_box(slide5, Inches(1.3), Inches(2.5), Inches(4.4), Inches(1.0),
             "重伤 2 人以上", font_size=36, bold=True, color=RED,
             alignment=PP_ALIGN.CENTER)
add_text_box(slide5, Inches(1.3), Inches(3.8), Inches(4.4), Inches(1.0),
             "达到此条件即启动应急预案", font_size=16, color=MED_GRAY,
             alignment=PP_ALIGN.CENTER)

add_shape_bg(slide5, Inches(7.3), Inches(2.0), Inches(5.0), Inches(3.5), LIGHT_RED)
add_text_box(slide5, Inches(7.6), Inches(2.5), Inches(4.4), Inches(1.0),
             "死亡 1 人以上", font_size=36, bold=True, color=RED,
             alignment=PP_ALIGN.CENTER)
add_text_box(slide5, Inches(7.6), Inches(3.8), Inches(4.4), Inches(1.0),
             "达到此条件即启动应急预案", font_size=16, color=MED_GRAY,
             alignment=PP_ALIGN.CENTER)


# ============================================
# SLIDE 6: Emergency Response Overview
# ============================================
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide6, WHITE)
add_section_header_bar(slide6, "二、应急响应过程 — 总览")
add_footer(slide6)

# Process flow boxes
steps = [
    ("① 接警与通知", "事故发生后\n立即报告校长\n启动应急预案"),
    ("② 现场抢救", "检查受伤情况\n应急救护处置\n送医院救治"),
    ("③ 联络与教育", "15分钟内\n写出书面报告\n稳定师生情绪"),
    ("④ 接待家长", "通知家长\n接待来访\n后勤支援保障"),
    ("⑤ 事故调查", "配合调查\n形成书面报告\n总结经验教训"),
]

box_width = Inches(2.1)
gap = Inches(0.25)
start_x = Inches(0.7)
y = Inches(2.0)

for i, (title, desc) in enumerate(steps):
    x = start_x + i * (box_width + gap)
    # Box
    box = add_shape_bg(slide6, x, y, box_width, Inches(2.8), LIGHT_GRAY)
    # Title
    add_text_box(slide6, x, y + Inches(0.2), box_width, Inches(0.5),
                 title, font_size=16, bold=True, color=RED,
                 alignment=PP_ALIGN.CENTER)
    # Description
    add_text_box(slide6, x + Inches(0.15), y + Inches(0.8), box_width - Inches(0.3), Inches(1.8),
                 desc, font_size=14, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)
    # Arrow between boxes
    if i < len(steps) - 1:
        arrow_x = x + box_width + Inches(0.02)
        add_text_box(slide6, arrow_x, y + Inches(1.0), Inches(0.2), Inches(0.5),
                     "→", font_size=24, bold=True, color=RED,
                     alignment=PP_ALIGN.CENTER)


# ============================================
# SLIDE 7: Alert & Notification
# ============================================
slide7 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide7, WHITE)
add_section_header_bar(slide7, "二、应急响应 — (一) 接警与通知")
add_footer(slide7)

add_shape_bg(slide7, Inches(0.8), Inches(1.3), Inches(11.7), Inches(5.3), LIGHT_GRAY)

items = [
    "事故发生后，在场人员必须立即制止，防止事态扩大，并报告校长",
    "校长必须掌握：事故发生的时间与地点、班级、程度",
    "基本掌握情况后，立即启动应急预案，值日行政和政教处领导赶赴现场",
    "上报教育局和医疗机构，通报包括：学校名称地址、通报人信息、事故时间地点种类程度、已采取的应急行动",
]
add_bullet_list(slide7, Inches(1.2), Inches(1.5), Inches(10.9), Inches(5.0),
                items, font_size=18, color=DARK_GRAY, spacing=Pt(14))


# ============================================
# SLIDE 8: On-site Rescue
# ============================================
slide8 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide8, WHITE)
add_section_header_bar(slide8, "二、应急响应 — (二) 现场应急抢救、现场保护")
add_footer(slide8)

add_shape_bg(slide8, Inches(0.8), Inches(1.3), Inches(11.7), Inches(5.3), LIGHT_GRAY)

items = [
    "在场人员首先检查学生受伤情况，先重后轻进行应急救护处置（或拨打120）",
    "校医到达后接替救护，确认需送医院救治的伤者，拨打120或派车送医",
    "门卫确保急救车进校后有人引导，组长向急救人员报告情况",
    "班主任及时通知家长简述事故情况和学生被送往的医院详细地址",
    "政教处组织调查事件过程，用分隔调查形式，实事求是做好书面记录",
    "严格保护事故现场，移动现场物件时必须做好标志、拍照、记录",
]
add_bullet_list(slide8, Inches(1.2), Inches(1.5), Inches(10.9), Inches(5.0),
                items, font_size=17, color=DARK_GRAY, spacing=Pt(12))


# ============================================
# SLIDE 9: Communication & Education
# ============================================
slide9 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide9, WHITE)
add_section_header_bar(slide9, "二、应急响应 — (三) 联络、教育")
add_footer(slide9)

add_shape_bg(slide9, Inches(0.8), Inches(1.3), Inches(11.7), Inches(5.3), LIGHT_GRAY)

items = [
    "办公室在15分钟内写出书面报告，经校长审查同意后送交教育局",
    "报告内容：时间地点、简要经过、伤亡人数、原因性质、抢救措施、协助需求",
    "做好师生教育工作，稳定情绪，严禁以个人名义向外扩散消息",
    "对情绪反应较大者安排心理教师进行辅导",
    "新闻媒体采访必须经过校长或上级部门同意，统一对外发布",
    "办公室参与事故调查处理，负责写出书面报告",
]
add_bullet_list(slide9, Inches(1.2), Inches(1.5), Inches(10.9), Inches(5.0),
                items, font_size=17, color=DARK_GRAY, spacing=Pt(12))


# ============================================
# SLIDE 10: Parent Reception & Logistics
# ============================================
slide10 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide10, WHITE)
add_section_header_bar(slide10, "二、应急响应 — (四) 接待家长和后勤支援")
add_footer(slide10)

add_shape_bg(slide10, Inches(0.8), Inches(1.3), Inches(11.7), Inches(5.3), LIGHT_GRAY)

items = [
    "看望、援助、救助伤亡学生家庭，做好家长思想工作和接待工作",
    "总务主任协助处理死亡学生的善后工作，安排住院学生家长的食宿",
    "依法调解安抚，掌握合法、合理、合情的原则，不留尾巴",
    "政教主任起草《协议书》，写明双方身份、事故经过、补偿协议",
    "整理病历卡复印件、医药费发票原件和复印件报保险公司理赔",
    "校安处组织保安严格核查外来人员身份，保证校园治安秩序稳定",
]
add_bullet_list(slide10, Inches(1.2), Inches(1.5), Inches(10.9), Inches(5.0),
                items, font_size=17, color=DARK_GRAY, spacing=Pt(12))


# ============================================
# SLIDE 11: Accident Investigation
# ============================================
slide11 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide11, WHITE)
add_section_header_bar(slide11, "二、应急响应 — (五) 事故调查")
add_footer(slide11)

add_shape_bg(slide11, Inches(0.8), Inches(1.3), Inches(11.7), Inches(5.3), LIGHT_GRAY)

items = [
    "配合上级部门进行事故处理及调查工作，调查事故原因",
    "整理事故记录，形成书面报告，向教育局报告事故处理结果",
    "对违反预案、不履行救援工作、发布假消息的人员进行处分",
    "构成犯罪的，移送司法机关依法追究刑事责任",
    "总结经验教训，查找制度、管理等存在的问题，制定防范措施",
]
add_bullet_list(slide11, Inches(1.2), Inches(1.5), Inches(10.9), Inches(5.0),
                items, font_size=18, color=DARK_GRAY, spacing=Pt(14))


# ============================================
# SLIDE 12: Summary / Closing
# ============================================
slide12 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide12, DARK_BG)
add_shape_bg(slide12, Inches(0), Inches(0), Inches(13.333), Inches(0.15), RED)

add_text_box(slide12, Inches(1.5), Inches(2.0), Inches(10), Inches(1.0),
             "安全第一 预防为主", font_size=40, bold=True, color=WHITE,
             alignment=PP_ALIGN.CENTER, font_name='微软雅黑')

add_text_box(slide12, Inches(1.5), Inches(3.3), Inches(10), Inches(2.0),
             "统一指挥 · 分级响应\n"
             "岗位责任 · 互相配合\n"
             "最大限度减少损失\n"
             "保障校园安全稳定",
             font_size=22, color=RGBColor(0xCC, 0xCC, 0xCC),
             alignment=PP_ALIGN.CENTER, font_name='微软雅黑')

add_shape_bg(slide12, Inches(0), Inches(7.35), Inches(13.333), Inches(0.15), RED)


# Save
output_path = "/home/admin/.openclaw/workspace/嘉善信息技术工程学校_校园学生突发事件应急预案.pptx"
prs.save(output_path)
print(f"✅ PPT saved to: {output_path}")
