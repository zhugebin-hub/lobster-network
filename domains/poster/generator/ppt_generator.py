#!/usr/bin/env python3
"""
HTML+Playwright PPT Generator
=============================
将每页PPT幻灯片作为HTML页面设计，用Playwright渲染为高清图片，
最后用python-pptx组装为标准PPTX文件。

用法:
  python ppt_generator.py --slides-dir ./my_slides --output presentation.pptx

作者: qoder小龙虾
日期: 2026-06-20
"""

import asyncio
import base64
import os
import sys
import glob
import argparse
from pathlib import Path

# ============================================================================
# Part 1: 幻灯片 HTML 模板库
# ============================================================================

SLIDE_CSS_BASE = """
:root {
  --slide-w: 1280px;
  --slide-h: 720px;
  --bg-primary: #0a0a2e;
  --bg-secondary: #1a0a3e;
  --neon-cyan: #00d4ff;
  --neon-orange: #ff6b35;
  --neon-purple: #7b2ff7;
  --gold: #ffd700;
  --text-primary: #ffffff;
  --text-secondary: rgba(255,255,255,0.8);
  --text-muted: rgba(255,255,255,0.5);
  --font: 'PingFang SC', 'Noto Sans SC', 'Microsoft YaHei', sans-serif;
}
* { margin: 0; padding: 0; box-sizing: border-box; }
body {
  width: var(--slide-w);
  height: var(--slide-h);
  background: linear-gradient(135deg, #0a0a2e 0%, #1a0a3e 40%, #0d1b3e 70%, #0a0a2e 100%);
  font-family: var(--font);
  color: var(--text-primary);
  overflow: hidden;
}
.page {
  width: var(--slide-w);
  height: var(--slide-h);
  position: relative;
  overflow: hidden;
}
.page-number {
  position: absolute;
  bottom: 20px;
  right: 40px;
  font-size: 14px;
  color: rgba(255,255,255,0.25);
  font-weight: 300;
}
.accent-line {
  width: 80px;
  height: 4px;
  background: linear-gradient(90deg, var(--neon-cyan), var(--neon-orange));
  border-radius: 2px;
}
.gradient-text {
  background: linear-gradient(90deg, var(--neon-cyan), var(--neon-orange));
  -webkit-background-clip: text;
  -webkit-text-fill-color: transparent;
  background-clip: text;
}
.glass-card {
  background: rgba(255,255,255,0.04);
  border: 1px solid rgba(0,212,255,0.15);
  border-radius: 16px;
  padding: 32px;
}
.glow-border {
  border: 1px solid rgba(0,212,255,0.3);
  box-shadow: 0 0 15px rgba(0,212,255,0.2), inset 0 0 15px rgba(0,212,255,0.05);
  border-radius: 16px;
}
"""

def make_slide_html(content_html, page_num=None, extra_css=""):
    """生成完整的幻灯片 HTML"""
    page_number_html = f'<span class="page-number">{page_num:02d}</span>' if page_num else ''
    return f"""<!DOCTYPE html>
<html lang="zh-CN">
<head>
  <meta charset="UTF-8">
  <style>{SLIDE_CSS_BASE}\n{extra_css}</style>
</head>
<body>
<div class="page">
  {content_html}
  {page_number_html}
</div>
</body>
</html>"""


# ---- 预置幻灯片模板 ----

def title_slide(title, subtitle="", author="", date_str=""):
    """封面页"""
    return make_slide_html(f"""
  <div style="display:flex; flex-direction:column; justify-content:center; align-items:center;
              height:100%; text-align:center; padding:60px 100px;">
    <div class="accent-line" style="margin-bottom:32px;"></div>
    <h1 class="gradient-text" style="font-size:52px; font-weight:800; line-height:1.3;
        margin-bottom:20px;">{title}</h1>
    {"<p style='font-size:24px; color:var(--text-secondary); margin-bottom:40px;'>" + subtitle + "</p>" if subtitle else ""}
    {"<p style='font-size:18px; color:var(--text-muted);'>" + author + "</p>" if author else ""}
    {"<p style='font-size:16px; color:var(--text-muted); margin-top:8px;'>" + date_str + "</p>" if date_str else ""}
  </div>
""")


def content_slide(title, bullets, page_num=None):
    """文字要点页"""
    bullet_html = "".join(
        f'<li style="font-size:22px; color:var(--text-secondary); margin-bottom:20px; '
        f'padding-left:16px; list-style:none; position:relative;">'
        f'<span style="position:absolute; left:-20px; top:6px; width:8px; height:8px; '
        f'background:var(--neon-cyan); border-radius:50%; box-shadow:0 0 10px var(--neon-cyan);"></span>'
        f'{b}</li>'
        for b in bullets
    )
    return make_slide_html(f"""
  <div style="padding:60px 80px; display:flex; flex-direction:column; height:100%;">
    <div class="accent-line" style="margin-bottom:20px;"></div>
    <h2 style="font-size:36px; font-weight:700; margin-bottom:32px;">{title}</h2>
    <ul style="flex:1; display:flex; flex-direction:column; justify-content:center; padding-left:36px;">
      {bullet_html}
    </ul>
  </div>
""", page_num)


def two_column_slide(title, left_title, left_items, right_title, right_items, page_num=None):
    """双栏对比页"""
    def col_html(title, items):
        lis = "".join(f'<li style="font-size:18px; color:var(--text-secondary); margin-bottom:12px;">{i}</li>' for i in items)
        return f"""<div class="glass-card" style="height:100%;">
          <h3 style="font-size:24px; color:var(--neon-cyan); margin-bottom:20px;">{title}</h3>
          <ul style="padding-left:20px;">{lis}</ul>
        </div>"""
    
    return make_slide_html(f"""
  <div style="padding:60px 80px; display:flex; flex-direction:column; height:100%;">
    <div class="accent-line" style="margin-bottom:20px;"></div>
    <h2 style="font-size:36px; font-weight:700; margin-bottom:32px;">{title}</h2>
    <div style="flex:1; display:grid; grid-template-columns:1fr 1fr; gap:32px; align-items:stretch;">
      {col_html(left_title, left_items)}
      {col_html(right_title, right_items)}
    </div>
  </div>
""", page_num)


def card_grid_slide(title, cards, page_num=None):
    """卡片网格页（2-6个卡片）"""
    cols = min(len(cards), 3)
    rows = (len(cards) + cols - 1) // cols
    
    card_html = ""
    for icon, card_title, desc in cards:
        card_html += f"""<div class="glass-card" style="text-align:center;">
          <div style="font-size:36px; margin-bottom:12px;">{icon}</div>
          <h4 style="font-size:20px; margin-bottom:8px; color:var(--text-primary);">{card_title}</h4>
          <p style="font-size:16px; color:var(--text-muted); line-height:1.6;">{desc}</p>
        </div>"""
    
    return make_slide_html(f"""
  <div style="padding:60px 80px; display:flex; flex-direction:column; height:100%;">
    <div class="accent-line" style="margin-bottom:20px;"></div>
    <h2 style="font-size:36px; font-weight:700; margin-bottom:32px;">{title}</h2>
    <div style="flex:1; display:grid; grid-template-columns:repeat({cols},1fr); 
                grid-template-rows:repeat({rows},1fr); gap:24px; align-items:stretch;">
      {card_html}
    </div>
  </div>
""", page_num)


def image_text_slide(title, text_content, image_path, page_num=None, image_left=True):
    """图文混排页"""
    img_tag = f'<img src="{image_path}" style="width:100%; height:100%; object-fit:cover; border-radius:12px;" />'
    text_block = f"""<div style="display:flex; flex-direction:column; justify-content:center;">
      <div class="accent-line" style="margin-bottom:20px;"></div>
      <h2 style="font-size:32px; font-weight:700; margin-bottom:20px;">{title}</h2>
      <p style="font-size:20px; color:var(--text-secondary); line-height:1.8;">{text_content}</p>
    </div>"""
    
    if image_left:
        left, right = img_tag, text_block
    else:
        left, right = text_block, img_tag
    
    return make_slide_html(f"""
  <div style="padding:60px 80px; display:flex; flex-direction:column; height:100%;">
    <div style="flex:1; display:grid; grid-template-columns:1fr 1fr; gap:48px; align-items:center;">
      {left}
      {right}
    </div>
  </div>
""", page_num)


def stats_slide(title, stats, page_num=None):
    """数据统计页 — stats: [(number, label), ...]"""
    items = ""
    for num, label in stats:
        items += f"""<div style="text-align:center;">
          <div class="gradient-text" style="font-size:56px; font-weight:800;">{num}</div>
          <p style="font-size:18px; color:var(--text-muted); margin-top:8px;">{label}</p>
        </div>"""
    
    return make_slide_html(f"""
  <div style="padding:60px 80px; display:flex; flex-direction:column; height:100%;">
    <div class="accent-line" style="margin-bottom:20px;"></div>
    <h2 style="font-size:36px; font-weight:700; margin-bottom:40px;">{title}</h2>
    <div style="flex:1; display:flex; justify-content:space-around; align-items:center;">
      {items}
    </div>
  </div>
""", page_num)


def thank_you_slide(title="谢谢", subtitle="", contact=""):
    """结束致谢页"""
    return make_slide_html(f"""
  <div style="display:flex; flex-direction:column; justify-content:center; align-items:center;
              height:100%; text-align:center; padding:60px 100px;">
    <h1 style="font-size:64px; font-weight:800; color:var(--neon-cyan);
        text-shadow: 0 0 20px rgba(0,212,255,0.5), 0 0 40px rgba(0,212,255,0.2);
        margin-bottom:24px;">{title}</h1>
    {"<p style='font-size:24px; color:var(--text-secondary); margin-bottom:32px;'>" + subtitle + "</p>" if subtitle else ""}
    {"<p style='font-size:18px; color:var(--text-muted);'>" + contact + "</p>" if contact else ""}
    <div class="accent-line" style="margin-top:40px; width:120px;"></div>
  </div>
""")


# ============================================================================
# Part 2: 渲染引擎
# ============================================================================

async def render_slides(slides_dir, output_dir, scale=2):
    """批量渲染目录下所有 slide_*.html 为 PNG"""
    from playwright.async_api import async_playwright
    
    os.makedirs(output_dir, exist_ok=True)
    html_files = sorted(glob.glob(os.path.join(slides_dir, 'slide_*.html')))
    
    if not html_files:
        print(f'No slide_*.html files found in {slides_dir}')
        return []
    
    print(f'Rendering {len(html_files)} slides...')
    
    async with async_playwright() as p:
        browser = await p.chromium.launch(headless=True)
        page = await browser.new_page(
            viewport={'width': 1280, 'height': 720},
            device_scale_factor=scale
        )
        
        png_files = []
        for html_file in html_files:
            await page.goto(f'file://{html_file}', wait_until='networkidle')
            await page.wait_for_timeout(1500)
            
            name = os.path.splitext(os.path.basename(html_file))[0]
            png_path = os.path.join(output_dir, f'{name}.png')
            await page.locator('.page').screenshot(path=png_path, type='png')
            png_files.append(png_path)
            print(f'  ✓ {name}.png')
        
        await browser.close()
    
    print(f'All {len(png_files)} slides rendered.')
    return png_files


# ============================================================================
# Part 3: PPTX 组装
# ============================================================================

def assemble_pptx(png_files, output_path):
    """将幻灯片 PNG 图片组装为 PPTX"""
    from pptx import Presentation
    from pptx.util import Inches, Emu
    
    SLIDE_WIDTH = Inches(13.333)
    SLIDE_HEIGHT = Inches(7.5)
    
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    
    blank_layout = prs.slide_layouts[6]  # 空白布局
    
    for png_path in png_files:
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(
            png_path,
            left=Emu(0), top=Emu(0),
            width=SLIDE_WIDTH, height=SLIDE_HEIGHT
        )
    
    prs.save(output_path)
    print(f'✓ PPTX saved: {output_path} ({len(png_files)} slides)')


# ============================================================================
# Part 4: 完整流程
# ============================================================================

async def generate_ppt(slides_dir, output_dir, output_pptx, scale=2):
    """完整的 PPT 生成流程"""
    print('=' * 50)
    print('  HTML+Playwright PPT Generator')
    print('=' * 50)
    
    # Step 1: 渲染
    print('\n📸 Step 1: Rendering slides to PNG...')
    pngs = await render_slides(slides_dir, output_dir, scale)
    
    if not pngs:
        print('❌ No slides rendered.')
        return
    
    # Step 2: 组装
    print('\n📦 Step 2: Assembling PPTX...')
    assemble_pptx(pngs, output_pptx)
    
    print(f'\n✅ Done! {len(pngs)} slides → {output_pptx}')


# ============================================================================
# CLI
# ============================================================================

if __name__ == '__main__':
    parser = argparse.ArgumentParser(description='HTML+Playwright PPT Generator')
    parser.add_argument('--slides-dir', required=True, help='Directory containing slide_*.html files')
    parser.add_argument('--output-dir', default='./slides_png', help='Output directory for PNG files')
    parser.add_argument('--output', default='presentation.pptx', help='Output PPTX file path')
    parser.add_argument('--scale', type=int, default=2, help='Render scale (default: 2 for Retina)')
    args = parser.parse_args()
    
    asyncio.run(generate_ppt(args.slides_dir, args.output_dir, args.output, args.scale))
