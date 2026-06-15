#!/usr/bin/env python3
"""生成 QPPB技术详解 PPT 文件（.pptx 格式）"""

from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

prs = Presentation()
prs.slide_width = Cm(33.866)  # 16:9
prs.slide_height = Cm(19.05)

# === 颜色定义 ===
DARK_BLUE = RGBColor(0, 51, 102)
MEDIUM_BLUE = RGBColor(0, 102, 153)
LIGHT_BLUE = RGBColor(200, 230, 255)
WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(0, 0, 0)
GRAY = RGBColor(128, 128, 128)
LIGHT_GRAY = RGBColor(240, 240, 240)
RED = RGBColor(204, 0, 0)

# === 中文字体设置 helper ===
def set_font(run, name='微软雅黑', size=Pt(18), color=BLACK, bold=False):
    run.font.name = name
    run.font.size = size
    run.font.color.rgb = color
    run.font.bold = bold
    # 设置中文字体
    rPr = run._element.get_or_add_rPr()
    rFonts = rPr.find(qn('a:rFonts'))
    if rFonts is None:
        rFonts = qn('a:rFonts')
        rPr.set(rFonts, name)
    # 设置 eastAsia 字体
    ea = rPr.find(qn('a:latin'))
    if ea is not None:
        ea.set('typeface', name)
    ea2 = rPr.find(qn('a:ea'))
    if ea2 is not None:
        ea2.set('typeface', name)

def add_shape_placeholder(slide, left, top, width, height, text='📷 图片占位', color=LIGHT_GRAY):
    """添加图片占位框"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = RGBColor(180, 180, 180)
    shape.line.width = Pt(1.5)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    set_font(run, name='微软雅黑', size=Pt(14), color=GRAY, bold=False)
    run.text = text
    return shape

def add_title_bar(slide, title_text):
    """添加顶部标题栏"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Cm(0), Cm(0), Cm(33.866), Cm(2.2)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    set_font(run, name='微软雅黑', size=Pt(24), color=WHITE, bold=True)
    run.text = title_text
    # 添加底部装饰线
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Cm(0), Cm(2.2), Cm(33.866), Cm(0.08)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = MEDIUM_BLUE
    line.line.fill.background()

def add_slide_number(slide, num, total=15):
    """添加页码"""
    txBox = slide.shapes.add_textbox(Cm(30), Cm(18.2), Cm(3.5), Cm(0.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    set_font(run, name='微软雅黑', size=Pt(10), color=GRAY)
    run.text = f'{num} / {total}'

def add_footer(slide):
    """添加底部页脚"""
    txBox = slide.shapes.add_textbox(Cm(1), Cm(18.2), Cm(10), Cm(0.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    set_font(run, name='微软雅黑', size=Pt(9), color=GRAY)
    run.text = 'QPPB技术详解 | 芦熠檑 | 2026'


# ============================================================
# Slide 1: 封面
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

# 背景色
bg = slide.background
fill = bg.fill
fill.solid()
fill.fore_color.rgb = WHITE

# 顶部装饰条
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(0), Cm(0), Cm(33.866), Cm(0.3))
shape.fill.solid()
shape.fill.fore_color.rgb = DARK_BLUE
shape.line.fill.background()

# 主标题
txBox = slide.shapes.add_textbox(Cm(3), Cm(5), Cm(28), Cm(3))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
set_font(run, name='微软雅黑', size=Pt(40), color=DARK_BLUE, bold=True)
run.text = 'QPPB技术详解'

# 副标题
txBox = slide.shapes.add_textbox(Cm(3), Cm(8.5), Cm(28), Cm(1.5))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
set_font(run, name='微软雅黑', size=Pt(20), color=MEDIUM_BLUE, bold=False)
run.text = '—— 基于 BGP 路由属性的 QoS 策略传播机制 ——'

# 分隔线
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(12), Cm(10.5), Cm(10), Cm(0.05))
line.fill.solid()
line.fill.fore_color.rgb = MEDIUM_BLUE
line.line.fill.background()

# 汇报人信息
txBox = slide.shapes.add_textbox(Cm(8), Cm(12), Cm(18), Cm(4))
tf = txBox.text_frame
for text in ['汇报人：芦熠檑', '日 期：2026 年 5 月']:
    p = tf.add_paragraph()
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    set_font(run, name='微软雅黑', size=Pt(16), color=GRAY, bold=False)
    run.text = text

prs.save('/home/admin/.openclaw/workspace/QPPB技术详解.pptx')
print('✅ Slide 1 封面完成')


# ============================================================
# Slide 2: 目录
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, '目  录')
add_slide_number(slide, 2)
add_footer(slide)

toc_items = [
    ('01', 'QPPB 概述与产生背景'),
    ('02', 'QPPB 实现原理'),
    ('03', 'QPPB 典型应用场景'),
]

for i, (num, title) in enumerate(toc_items):
    y = Cm(4 + i * 3.5)
    # 序号
    txBox = slide.shapes.add_textbox(Cm(4), y, Cm(3), Cm(2))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    set_font(run, name='微软雅黑', size=Pt(36), color=DARK_BLUE, bold=True)
    run.text = num

    # 标题
    txBox = slide.shapes.add_textbox(Cm(7.5), y + Cm(0.3), Cm(20), Cm(1.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    set_font(run, name='微软雅黑', size=Pt(22), color=BLACK, bold=False)
    run.text = title

    # 装饰线
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(4), y + Cm(2.5), Cm(26), Cm(0.03))
    line.fill.solid()
    line.fill.fore_color.rgb = LIGHT_BLUE
    line.line.fill.background()

prs.save('/home/admin/.openclaw/workspace/QPPB技术详解.pptx')
print('✅ Slide 2 目录完成')


# ============================================================
# Slide 3: 第一部分标题页
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
# 背景色
bg = slide.background
fill = bg.fill
fill.solid()
fill.fore_color.rgb = DARK_BLUE

txBox = slide.shapes.add_textbox(Cm(4), Cm(7), Cm(26), Cm(3))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
set_font(run, name='微软雅黑', size=Pt(36), color=WHITE, bold=True)
run.text = '第一部分'

txBox = slide.shapes.add_textbox(Cm(4), Cm(10.5), Cm(26), Cm(2))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
set_font(run, name='微软雅黑', size=Pt(24), color=LIGHT_BLUE, bold=False)
run.text = 'QPPB 概述与产生背景'

add_slide_number(slide, 3)

prs.save('/home/admin/.openclaw/workspace/QPPB技术详解.pptx')
print('✅ Slide 3 第一部分标题完成')


# ============================================================
# Slide 4: 什么是 QPPB？
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, '1.1 什么是 QPPB？')
add_slide_number(slide, 4)
add_footer(slide)

items = [
    ('全称', 'QoS Policy Propagation through BGP\n（通过 BGP 传播 QoS 策略）'),
    ('定义', '一种特殊的复杂流分类方法，通过 BGP 路由属性对报文进行流分类'),
    ('核心思想', 'BGP 路由发送者通过设置 BGP 属性预先对路由进行分类，\nBGP 路由接收者根据路由属性匹配并关联 QoS 策略'),
    ('优势', '网络结构变化时，只需修改发送端配置，\n接收端无需改动'),
]

for i, (label, content) in enumerate(items):
    y = Cm(3 + i * 3.5)
    # 标签
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(2), y, Cm(3), Cm(1.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    set_font(run, name='微软雅黑', size=Pt(14), color=WHITE, bold=True)
    run.text = label

    # 内容
    txBox = slide.shapes.add_textbox(Cm(5.5), y + Cm(0.1), Cm(25), Cm(1.2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    set_font(run, name='微软雅黑', size=Pt(14), color=BLACK, bold=False)
    run.text = content

prs.save('/home/admin/.openclaw/workspace/QPPB技术详解.pptx')
print('✅ Slide 4 什么是 QPPB 完成')


# ============================================================
# Slide 5: 产生背景（痛点分析）
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, '1.2 产生背景（痛点分析）')
add_slide_number(slide, 5)
add_footer(slide)

# 左侧：痛点
txBox = slide.shapes.add_textbox(Cm(2), Cm(3), Cm(14), Cm(1.2))
tf = txBox.text_frame
p = tf.paragraphs[0]
run = p.add_run()
set_font(run, name='微软雅黑', size=Pt(16), color=RED, bold=True)
run.text = '❌ 传统方式的痛点'

problems = [
    'AS400 是高优先级网络，需对往返报文重新设置 IP Precedence',
    'Node-A/Node-B 需针对 AS400 内大量 IP 地址/地址段配置流分类',
    '网络结构不稳定时，配置修改工作量巨大',
    '大量流分类规则难以维护',
]

for i, prob in enumerate(problems):
    y = Cm(4.8 + i * 1.8)
    txBox = slide.shapes.add_textbox(Cm(2.5), y, Cm(13), Cm(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    set_font(run, name='微软雅黑', size=Pt(12), color=BLACK, bold=False)
    run.text = f'• {prob}'

# 右侧：解决方案
txBox = slide.shapes.add_textbox(Cm(18), Cm(3), Cm(14), Cm(1.2))
tf = txBox.text_frame
p = tf.paragraphs[0]
run = p.add_run()
set_font(run, name='微软雅黑', size=Pt(16), color=DARK_BLUE, bold=True)
run.text = '✅ QPPB 解决方案'

solutions = [
    '按 AS 信息、团体属性等聚类信息对报文分类',
    '发送端设置路由属性，接收端自动匹配',
    '网络变化只需修改发送端配置',
    '大幅简化配置和维护工作量',
]

for i, sol in enumerate(solutions):
    y = Cm(4.8 + i * 1.8)
    txBox = slide.shapes.add_textbox(Cm(18.5), y, Cm(13), Cm(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    set_font(run, name='微软雅黑', size=Pt(12), color=BLACK, bold=False)
    run.text = f'✓ {sol}'

# 分隔线
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(16.5), Cm(3), Cm(0.03), Cm(12))
line.fill.solid()
line.fill.fore_color.rgb = LIGHT_BLUE
line.line.fill.background()

# 图片占位
add_shape_placeholder(slide, Cm(2), Cm(13), Cm(30), Cm(5), '📷 图 1-52 跨 AS 组网示意图')

prs.save('/home/admin/.openclaw/workspace/QPPB技术详解.pptx')
print('✅ Slide 5 产生背景完成')


# ============================================================
# Slide 6: QPPB技术定位
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, '1.3 QPPB技术定位')
add_slide_number(slide, 6)
add_footer(slide)

# 对比表格
from pptx.util import Inches, Pt, Emu
from pptx.oxml.xmlchemy import OxmlElement

table_data = [
    ['对比维度', '简单流分类', 'QPPB（复杂流分类）'],
    ['分类依据', '报文头部固定字段\n（DSCP / 802.1p）', 'BGP 路由属性\n（AS_PATH / Community）'],
    ['配置复杂度', '低（直接基于报文头部）', '中（需配置 BGP 属性）'],
    ['适用场景', '单域/简单网络', '跨 AS / 大型复杂组网'],
    ['灵活性', '低（网络变化需重新配置）', '高（发送端修改即可）'],
    ['维护成本', '高（需逐设备配置）', '低（接收端自动适配）'],
]

rows, cols = len(table_data), len(table_data[0])
table = slide.shapes.add_table(rows, cols, Cm(2), Cm(3), Cm(30), Cm(12)).table

# 设置列宽
table.columns[0].width = Cm(5)
table.columns[1].width = Cm(12.5)
table.columns[2].width = Cm(12.5)

for row_idx, row_data in enumerate(table_data):
    for col_idx, cell_text in enumerate(row_data):
        cell = table.cell(row_idx, col_idx)
        cell.text = ''
        p = cell.text_frame.paragraphs[0]
        run = p.add_run()
        set_font(run, name='微软雅黑', size=Pt(12), color=BLACK, bold=False)
        run.text = cell_text
        p.alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

        # 表头样式
        if row_idx == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = DARK_BLUE
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    run.font.color.rgb = WHITE
                    run.font.bold = True
        elif row_idx % 2 == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT_GRAY

prs.save('/home/admin/.openclaw/workspace/QPPB技术详解.pptx')
print('✅ Slide 6 技术定位完成')


# ============================================================
# Slide 7: 第二部分标题页
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.background
fill = bg.fill
fill.solid()
fill.fore_color.rgb = DARK_BLUE

txBox = slide.shapes.add_textbox(Cm(4), Cm(7), Cm(26), Cm(3))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
set_font(run, name='微软雅黑', size=Pt(36), color=WHITE, bold=True)
run.text = '第二部分'

txBox = slide.shapes.add_textbox(Cm(4), Cm(10.5), Cm(26), Cm(2))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
set_font(run, name='微软雅黑', size=Pt(24), color=LIGHT_BLUE, bold=False)
run.text = 'QPPB 实现原理'

add_slide_number(slide, 7)

prs.save('/home/admin/.openclaw/workspace/QPPB技术详解.pptx')
print('✅ Slide 7 第二部分标题完成')


# ============================================================
# Slide 8: QPPB 工作流程
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, '2.1 QPPB 工作流程（图 1-53）')
add_slide_number(slide, 8)
add_footer(slide)

steps = [
    ('步骤 1', 'BGP 路由发送者\n（Node-C）为 BGP 路由\n设置特定属性\n（AS_PATH / Community）'),
    ('步骤 2', 'BGP 路由携带属性\n在 AS 间通告\n（属性作为分类标识）'),
    ('步骤 3', 'BGP 路由接收者\n（Node-A）匹配路由属性\n设置 Behavior ID 到 FIB'),
    ('步骤 4', '数据转发时\n根据目的网络从 FIB\n获取 Behavior ID\n执行对应流动作'),
]

for i, (label, content) in enumerate(steps):
    x = Cm(1.5 + i * 8)
    y = Cm(3)

    # 步骤框
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Cm(7), Cm(7))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = DARK_BLUE
    shape.line.width = Pt(2)

    # 标签
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    set_font(run, name='微软雅黑', size=Pt(14), color=DARK_BLUE, bold=True)
    run.text = label + '\n\n'

    # 内容
    p = tf.add_paragraph()
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    set_font(run, name='微软雅黑', size=Pt(11), color=BLACK, bold=False)
    run.text = content

    # 箭头（除了最后一个）
    if i < 3:
        arrow_x = x + Cm(7.2)
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, arrow_x, Cm(5.5), Cm(0.8), Cm(1))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = MEDIUM_BLUE
        arrow.line.fill.background()

# 图片占位
add_shape_placeholder(slide, Cm(2), Cm(11.5), Cm(30), Cm(6), '📷 图 1-53 QPPB 实现原理示意图')

prs.save('/home/admin/.openclaw/workspace/QPPB技术详解.pptx')
print('✅ Slide 8 工作流程完成')


# ============================================================
# Slide 9: 关键技术点
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, '2.2 关键技术点')
add_slide_number(slide, 9)
add_footer(slide)

key_points = [
    ('Behavior ID', '不同的流动作对应不同的 Behavior ID，存储在 FIB 表项中。\n数据转发时，设备根据目的网络从 FIB 获取 Behavior ID，\n执行相应的流动作。'),
    ('QoS Local-ID', 'QPPB 策略中绑定 qos-local-id 与 behavior，\n实现路由属性与 QoS 策略的关联。'),
    ('策略传递机制', '路由发送端：通过 route-policy 设置 AS_PATH / Community / Ext-Community\n路由接收端：通过 route-policy import 匹配属性，apply qos-local-id'),
]

for i, (title, content) in enumerate(key_points):
    y = Cm(3 + i * 4.5)

    # 标题框
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(2), y, Cm(6), Cm(1.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    set_font(run, name='微软雅黑', size=Pt(16), color=WHITE, bold=True)
    run.text = title

    # 内容
    txBox = slide.shapes.add_textbox(Cm(8.5), y + Cm(0.1), Cm(23), Cm(1.2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    set_font(run, name='微软雅黑', size=Pt(12), color=BLACK, bold=False)
    run.text = content

# 重要说明
txBox = slide.shapes.add_textbox(Cm(2), Cm(15.5), Cm(30), Cm(2.5))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
run = p.add_run()
set_font(run, name='微软雅黑', size=Pt(13), color=RED, bold=True)
run.text = '⚠️ 重要说明：'
p = tf.add_paragraph()
run = p.add_run()
set_font(run, name='微软雅黑', size=Pt(12), color=BLACK, bold=False)
run.text = 'QPPB技术实际并没有在 BGP 路由信息中发送 QoS 策略，\n只是在路由发送方通过对通告的路由设置路由属性，\n在路由接收方根据目的网段的路由属性设置 QoS 策略。'

prs.save('/home/admin/.openclaw/workspace/QPPB技术详解.pptx')
print('✅ Slide 9 关键技术点完成')


# ============================================================
# Slide 10: 上行 vs 下行 QPPB
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, '2.3 上行 vs 下行 QPPB')
add_slide_number(slide, 10)
add_footer(slide)

table_data = [
    ['方向', '配置命令', '查表依据', '应用场景'],
    ['上行\n（inbound）', 'qppb-policy policy\nsource inbound', '根据源 IP\n查路由表', '用户→ISP\n流量计费'],
    ['下行\n（outbound）', 'qppb-policy policy\noutbound', '根据目的 IP\n查路由表', 'ISP→用户\n流量计费'],
    ['基于 IP 优先级', 'qppb-policy ip-precedence\nsource', '根据源/目的地址', '按优先级分类'],
]

rows, cols = len(table_data), len(table_data[0])
table = slide.shapes.add_table(rows, cols, Cm(2), Cm(3), Cm(30), Cm(10)).table

table.columns[0].width = Cm(5)
table.columns[1].width = Cm(10)
table.columns[2].width = Cm(7)
table.columns[3].width = Cm(8)

for row_idx, row_data in enumerate(table_data):
    for col_idx, cell_text in enumerate(row_data):
        cell = table.cell(row_idx, col_idx)
        cell.text = ''
        p = cell.text_frame.paragraphs[0]
        run = p.add_run()
        set_font(run, name='微软雅黑', size=Pt(12), color=BLACK, bold=False)
        run.text = cell_text
        p.alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

        if row_idx == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = DARK_BLUE
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    run.font.color.rgb = WHITE
                    run.font.bold = True
        elif row_idx % 2 == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT_GRAY

prs.save('/home/admin/.openclaw/workspace/QPPB技术详解.pptx')
print('✅ Slide 10 上行 vs 下行完成')


# ============================================================
# Slide 11: 第三部分标题页
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.background
fill = bg.fill
fill.solid()
fill.fore_color.rgb = DARK_BLUE

txBox = slide.shapes.add_textbox(Cm(4), Cm(7), Cm(26), Cm(3))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
set_font(run, name='微软雅黑', size=Pt(36), color=WHITE, bold=True)
run.text = '第三部分'

txBox = slide.shapes.add_textbox(Cm(4), Cm(10.5), Cm(26), Cm(2))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
set_font(run, name='微软雅黑', size=Pt(24), color=LIGHT_BLUE, bold=False)
run.text = 'QPPB 典型应用场景'

add_slide_number(slide, 11)

prs.save('/home/admin/.openclaw/workspace/QPPB技术详解.pptx')
print('✅ Slide 11 第三部分标题完成')


# ============================================================
# Slide 12: 典型应用一 — AS 域间流量分类
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, '3.1 典型应用一：AS 域间流量分类（图 1-54）')
add_slide_number(slide, 12)
add_footer(slide)

# 场景描述
txBox = slide.shapes.add_textbox(Cm(2), Cm(3), Cm(30), Cm(1.5))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
run = p.add_run()
set_font(run, name='微软雅黑', size=Pt(14), color=BLACK, bold=False)
run.text = '场景：使用 QPPB 可以方便地在 AS100 的边缘设备对 AS 域间的流量进行流分类。\n例如要在 Node-C 上对 AS200 和 AS400 之间的流量进行限速。'

# 配置方案
txBox = slide.shapes.add_textbox(Cm(2), Cm(5), Cm(30), Cm(1))
tf = txBox.text_frame
p = tf.paragraphs[0]
run = p.add_run()
set_font(run, name='微软雅黑', size=Pt(16), color=DARK_BLUE, bold=True)
run.text = '配置方案：'

solutions = [
    'AS200 → AS400 方向：在 Node-C 上的 AS100 域内所有接口使能针对源地址的 QPPB',
    'AS400 → AS200 方向：在 Node-C 上与 AS400 相连的接口使能针对目的地址的 QPPB',
]

for i, sol in enumerate(solutions):
    y = Cm(6.5 + i * 1.5)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(2), y, Cm(30), Cm(1.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT_BLUE
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    set_font(run, name='微软雅黑', size=Pt(12), color=DARK_BLUE, bold=False)
    run.text = f'{i+1}. {sol}'

# 须知
txBox = slide.shapes.add_textbox(Cm(2), Cm(10.5), Cm(30), Cm(1.5))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
run = p.add_run()
set_font(run, name='微软雅黑', size=Pt(13), color=RED, bold=True)
run.text = '⚠️ 须知：查 FIB 转发的是针对上行流量而不是下行流量，因此使能 QPPB 的接口是流量上行的接口。'

# 图片占位
add_shape_placeholder(slide, Cm(2), Cm(12.5), Cm(30), Cm(5), '📷 图 1-54 AS 域间流量分类组网示意图')

prs.save('/home/admin/.openclaw/workspace/QPPB技术详解.pptx')
print('✅ Slide 12 应用一完成')


# ============================================================
# Slide 13: 典型应用二 — L3VPN 流量分类
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, '3.2 典型应用二：L3VPN 流量分类（图 1-55）')
add_slide_number(slide, 13)
add_footer(slide)

txBox = slide.shapes.add_textbox(Cm(2), Cm(3), Cm(30), Cm(1.5))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
run = p.add_run()
set_font(run, name='微软雅黑', size=Pt(14), color=BLACK, bold=False)
run.text = '场景：QPPB技术在 BGP/MPLS L3VPN 组网环境中的应用。\n当 PE 连接多个 VPN 时，可以对某个 VPN-instance 在路由发布时设置 Community 等属性后，再将路由通告出去。'

txBox = slide.shapes.add_textbox(Cm(2), Cm(5), Cm(30), Cm(1))
tf = txBox.text_frame
p = tf.paragraphs[0]
run = p.add_run()
set_font(run, name='微软雅黑', size=Pt(16), color=DARK_BLUE, bold=True)
run.text = '配置方案：'

steps = [
    'PE 连接多个 VPN 时，对某个 VPN-instance 在路由发布时设置 Community 等属性后，再将路由通告出去',
    '远端 PE 接收到路由信息后将路由及 QoS 等参数设置到 FIB 表项中',
    '使得从 CE 来的流量在转发时能执行相应的 QoS 动作',
    '这样，不同的 VPN 可获得不同的服务质量',
]

for i, step in enumerate(steps):
    y = Cm(6.5 + i * 1.5)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(2), y, Cm(30), Cm(1.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT_BLUE
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    set_font(run, name='微软雅黑', size=Pt(12), color=DARK_BLUE, bold=False)
    run.text = f'{i+1}. {step}'

# 图片占位
add_shape_placeholder(slide, Cm(2), Cm(13), Cm(30), Cm(5), '📷 图 1-55 L3VPN 流量分类组网示意图')

prs.save('/home/admin/.openclaw/workspace/QPPB技术详解.pptx')
print('✅ Slide 13 应用二完成')


# ============================================================
# Slide 14: 典型应用三 — 用户→ISP 流量计费
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, '3.3 典型应用三：用户→ISP 的流量计费（图 1-56）')
add_slide_number(slide, 14)
add_footer(slide)

txBox = slide.shapes.add_textbox(Cm(2), Cm(3), Cm(30), Cm(1))
tf = txBox.text_frame
p = tf.paragraphs[0]
run = p.add_run()
set_font(run, name='微软雅黑', size=Pt(14), color=BLACK, bold=False)
run.text = '场景：QPPB技术应用于用户到 ISP 的流量计费场景。'

txBox = slide.shapes.add_textbox(Cm(2), Cm(4.5), Cm(30), Cm(1))
tf = txBox.text_frame
p = tf.paragraphs[0]
run = p.add_run()
set_font(run, name='微软雅黑', size=Pt(16), color=DARK_BLUE, bold=True)
run.text = '配置方案：'

steps = [
    '通过 BGP 协议，发布路由时携带团体属性',
    '引入 BGP 路由时，匹配团体属性，在路由表中设置 Behavior ID',
    '配置 qppb-policy，匹配 qos-local-id，配置统计/CAR/Remark 等动作',
    '在流量入口方向使能基于目的地址的 QPPB',
    '在用户侧接口的 inbound 方向应用 qppb-policy',
]

for i, step in enumerate(steps):
    y = Cm(6 + i * 1.3)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(2), y, Cm(30), Cm(1.1))
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT_BLUE
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    set_font(run, name='微软雅黑', size=Pt(11), color=DARK_BLUE, bold=False)
    run.text = f'{i+1}. {step}'

# 转发过程
txBox = slide.shapes.add_textbox(Cm(2), Cm(13.5), Cm(30), C