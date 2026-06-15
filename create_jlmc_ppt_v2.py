#!/usr/bin/env python3
"""
吉林医药学院 PPT 模板 V2 - 完整版
包含：封面、目录、内容页(左文右图)、两列对比页、结束页
配色：深蓝 #00569E / 红色 #C41230 / 亮蓝 #0096D6
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# 配色方案
PRIMARY    = RGBColor(0x00, 0x56, 0x9E)   # 深蓝
SECONDARY  = RGBColor(0xC4, 0x12, 0x30)   # 红色
ACCENT     = RGBColor(0x00, 0x96, 0xD6)   # 亮蓝
LIGHT_BLUE = RGBColor(0xE8, 0xF4, 0xFC)   # 浅蓝背景
DARK       = RGBColor(0x33, 0x33, 0x33)   # 深色文字
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
GRAY       = RGBColor(0x99, 0x99, 0x99)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
MED_BLUE   = RGBColor(0x00, 0x6B, 0xB8)   # 中蓝（装饰用）
LIGHT_RED  = RGBColor(0xFF, 0xE8, 0xEB)   # 浅红背景

W = Inches  # 快捷引用

def add_shape(slide, shape_type, left, top, width, height, fill_color=None, line_color=None, line_width=None):
    """快捷添加形状"""
    s = slide.shapes.add_shape(shape_type, left, top, width, height)
    if fill_color:
        s.fill.solid()
        s.fill.fore_color.rgb = fill_color
    else:
        s.fill.background()
    if line_color:
        s.line.color.rgb = line_color
        s.line.width = Pt(line_width or 1)
    else:
        s.line.fill.background()
    return s

def add_textbox(slide, left, top, width, height, text, font_size=18, bold=False, color=DARK, alignment=PP_ALIGN.LEFT, word_wrap=True):
    """快捷添加文本框"""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    return tb

def set_slide_bg(slide, color):
    """设置幻灯片背景"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_page_number(slide, num, total=None):
    """添加页码"""
    text = f"{num}" if total is None else f"{num}/{total}"
    add_textbox(slide, W(11.5), W(7.0), W(1.5), W(0.4), text, 12, False, GRAY, PP_ALIGN.RIGHT)

def add_top_bar(slide):
    """顶部装饰条"""
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, W(13.333), W(0.06), PRIMARY)

def add_bottom_bar(slide):
    """底部装饰条"""
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0, W(7.35), W(13.333), W(0.15), SECONDARY)

def add_section_title(slide, title, y=W(0.5)):
    """添加章节标题+下划线"""
    add_textbox(slide, W(1), y, W(11.333), W(0.8), title, 32, True, PRIMARY)
    add_shape(slide, MSO_SHAPE.RECTANGLE, W(1), y + W(0.8), W(2.5), W(0.04), SECONDARY)

def create_ppt():
    prs = Presentation()
    prs.slide_width = W(13.333)
    prs.slide_height = W(7.5)
    total_slides = 5
    
    # ============ 第 1 页：封面 ============
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    
    # 顶部蓝色区域
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, W(13.333), W(2.8), PRIMARY)
    
    # 底部白色区域装饰线
    add_shape(slide, MSO_SHAPE.RECTANGLE, W(2), W(5.2), W(9.333), W(0.04), SECONDARY)
    
    # 学校名称（白色区域）
    add_textbox(slide, W(1.5), W(0.6), W(10.333), W(1.2), '吉林医药学院', 48, True, WHITE, PP_ALIGN.CENTER)
    
    # 英文校名
    add_textbox(slide, W(2), W(1.8), W(9.333), W(0.6), 'JILIN MEDICAL UNIVERSITY', 16, False, RGBColor(0xBB, 0xCC, 0xDD), PP_ALIGN.CENTER)
    
    # 主标题
    add_textbox(slide, W(1.5), W(3.3), W(10.333), W(1.0), '论文题目 / 报告主题', 36, True, DARK, PP_ALIGN.CENTER)
    
    # 汇报信息
    add_textbox(slide, W(2), W(4.6), W(9.333), W(0.8), 
                '汇报人：XXX    指导教师：XXX    日期：2026 年 X 月 X 日', 
                18, False, GRAY, PP_ALIGN.CENTER)
    
    # 底部红色条
    add_bottom_bar(slide)
    
    # ============ 第 2 页：目录 ============
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_top_bar(slide)
    add_section_title(slide, '目录')
    
    toc_items = [
        ('01', '研究背景与意义'),
        ('02', '文献综述'),
        ('03', '研究方法与内容'),
        ('04', '研究结果与分析'),
        ('05', '结论与展望'),
    ]
    
    for i, (num, text) in enumerate(toc_items):
        y = W(2.0 + i * 0.95)
        
        # 序号圆点
        circle = add_shape(slide, MSO_SHAPE.OVAL, W(1.3), y + W(0.05), W(0.5), W(0.5), SECONDARY)
        tf = circle.text_frame
        p = tf.paragraphs[0]
        p.text = num
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        
        # 目录文字
        add_textbox(slide, W(2.2), y, W(8), W(0.6), text, 24, False, DARK)
        
        # 分隔线（最后一条不画）
        if i < len(toc_items) - 1:
            add_shape(slide, MSO_SHAPE.RECTANGLE, W(1.5), y + W(0.75), W(9.5), W(0.01), GRAY)
    
    add_page_number(slide, 2, total_slides)
    
    # ============ 第 3 页：内容页（左文右图）============
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_top_bar(slide)
    add_section_title(slide, '研究背景与意义')
    
    # 左侧内容卡片
    left_card = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, W(0.6), W(1.8), W(5.8), W(5.2), LIGHT_BLUE, PRIMARY, 1.5)
    tf = left_card.text_frame
    tf.word_wrap = True
    
    left_content = [
        ('研究背景', True, 22, PRIMARY),
        ('', False, 8, DARK),
        ('• 随着我国医疗卫生事业的发展，', False, 16, DARK),
        ('  医学教育面临新的挑战和机遇', False, 16, DARK),
        ('', False, 8, DARK),
        ('• 新医科建设对医学人才培养', False, 16, DARK),
        ('  提出更高要求', False, 16, DARK),
        ('', False, 8, DARK),
        ('• 本研究立足于吉林医药学院实际，', False, 16, DARK),
        ('  探索创新路径', False, 16, DARK),
    ]
    
    for i, (text, bold, size, color) in enumerate(left_content):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.space_after = Pt(4)
    
    # 右侧图表占位区
    right_card = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, W(6.8), W(1.8), W(5.8), W(5.2), LIGHT_BLUE, ACCENT, 1.5)
    tf = right_card.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = '📊 图表 / 图片占位区\n\n建议尺寸：\n宽 5.5 英寸 × 高 4.8 英寸\n\n支持格式：PNG / JPG / SVG'
    p.font.size = Pt(18)
    p.font.color.rgb = PRIMARY
    p.alignment = PP_ALIGN.CENTER
    
    add_page_number(slide, 3, total_slides)
    
    # ============ 第 4 页：两列对比页 ============
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_top_bar(slide)
    add_section_title(slide, '文献综述')
    
    # 左列 - 国内
    left_col = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, W(0.6), W(1.8), W(5.8), W(5.0), LIGHT_BLUE, PRIMARY, 1.5)
    tf = left_col.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = '🇨🇳 国内研究现状'
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = PRIMARY
    
    p = tf.add_paragraph()
    p.text = '\n• 国内学者在医学教育领域的研究成果\n• 主要研究方法和理论框架\n• 存在的问题与不足\n• 未来发展趋势'
    p.font.size = Pt(16)
    p.font.color.rgb = DARK
    p.space_after = Pt(8)
    
    # 右列 - 国外
    right_col = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, W(6.8), W(1.8), W(5.8), W(5.0), LIGHT_GRAY, GRAY, 1)
    tf = right_col.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = '🌍 国外研究现状'
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = PRIMARY
    
    p = tf.add_paragraph()
    p.text = '\n• 国际医学教育发展趋势\n• 先进教育理念与模式\n• 可借鉴的经验与启示\n• 本土化应用建议'
    p.font.size = Pt(16)
    p.font.color.rgb = DARK
    p.space_after = Pt(8)
    
    add_page_number(slide, 4, total_slides)
    
    # ============ 第 5 页：结束页 ============
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, PRIMARY)
    
    # 装饰圆
    add_shape(slide, MSO_SHAPE.OVAL, W(4.2), W(0.8), W(4.933), W(4.933), MED_BLUE)
    
    # 感谢文字
    add_textbox(slide, W(2.5), W(2.5), W(8.333), W(1.5), '感谢聆听', 48, True, WHITE, PP_ALIGN.CENTER)
    
    # 学校名称
    add_textbox(slide, W(3), W(4.3), W(7.333), W(0.8), '吉林医药学院', 24, False, RGBColor(0xBB, 0xCC, 0xDD), PP_ALIGN.CENTER)
    
    # 底部装饰线
    add_shape(slide, MSO_SHAPE.RECTANGLE, W(3.5), W(5.8), W(6.333), W(0.04), SECONDARY)
    
    # 保存
    output = '/home/admin/.openclaw/workspace/吉林医药学院_PPT模板_完整版.pptx'
    prs.save(output)
    print(f'✅ PPT 模板已生成：{output}')
    return output

if __name__ == '__main__':
    create_ppt()
