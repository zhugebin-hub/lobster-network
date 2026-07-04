#!/usr/bin/env python3
"""生成《生成式人工智能技术应用》分模块PPT - 每个模块独立PPT，含课堂练习"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# 科技风格配色
BG_DARK = RGBColor(0x0A, 0x0E, 0x24)
BG_CARD = RGBColor(0x12, 0x1A, 0x3A)
ACCENT_CYAN = RGBColor(0x00, 0xD4, 0xFF)
ACCENT_BLUE = RGBColor(0x3B, 0x82, 0xF7)
ACCENT_PURPLE = RGBColor(0x8B, 0x5C, 0xF6)
ACCENT_GREEN = RGBColor(0x10, 0xB9, 0x81)
ACCENT_ORANGE = RGBColor(0xF5, 0x9E, 0x0B)
TEXT_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_LIGHT = RGBColor(0xC9, 0xD1, 0xD9)
TEXT_DIM = RGBColor(0x8B, 0x94, 0x9E)

def new_ppt():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs

def add_bg(slide):
    bg = slide.background; fill = bg.fill; fill.solid(); fill.fore_color.rgb = BG_DARK

def add_bar(slide, left=0, top=0, w=13.333, h=0.06, color=ACCENT_CYAN):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()

def add_tb(slide, text, left=0.5, top=0.3, w=12, h=0.8, size=32, color=ACCENT_CYAN, bold=True):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(w), Inches(h))
    tf = txBox.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.font.size = Pt(size); p.font.bold = bold; p.font.color.rgb = color
    return txBox

def add_body(slide, text, left=0.8, top=1.5, w=11.5, h=5.5, size=18, color=TEXT_LIGHT, bold=False):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(w), Inches(h))
    tf = txBox.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.font.size = Pt(size); p.font.color.rgb = color; p.font.bold = bold
    return txBox

def add_bullets(slide, items, left=0.8, top=1.5, w=11.5, h=5.5, size=18, color=TEXT_LIGHT, spacing=Pt(8)):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(w), Inches(h))
    tf = txBox.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item; p.font.size = Pt(size); p.font.color.rgb = color; p.space_after = spacing
    return txBox

def add_code(slide, code, left=0.8, top=1.5, w=11.5, h=2.5, size=16):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(w), Inches(h))
    shape.fill.solid(); shape.fill.fore_color.rgb = BG_CARD; shape.line.fill.background()
    txBox = slide.shapes.add_textbox(Inches(left+0.2), Inches(top+0.15), Inches(w-0.4), Inches(h-0.3))
    tf = txBox.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = code; p.font.size = Pt(size); p.font.color.rgb = ACCENT_GREEN; p.font.name = "Consolas"
    return txBox

def add_card(slide, title, content, left=0.8, top=1.5, w=5.5, h=2.5, tc=ACCENT_CYAN):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(w), Inches(h))
    shape.fill.solid(); shape.fill.fore_color.rgb = BG_CARD
    shape.line.color.rgb = tc; shape.line.width = Pt(1.5)
    txBox = slide.shapes.add_textbox(Inches(left+0.2), Inches(top+0.15), Inches(w-0.4), Inches(0.5))
    tf = txBox.text_frame; p = tf.paragraphs[0]; p.text = title; p.font.size = Pt(20); p.font.bold = True; p.font.color.rgb = tc
    txBox2 = slide.shapes.add_textbox(Inches(left+0.2), Inches(top+0.7), Inches(w-0.4), Inches(h-0.9))
    tf2 = txBox2.text_frame; tf2.word_wrap = True
    p2 = tf2.paragraphs[0]; p2.text = content; p2.font.size = Pt(15); p2.font.color.rgb = TEXT_LIGHT

def add_table(slide, headers, rows, left=0.8, top=1.5, w=11.5, h=4):
    cols = len(headers)
    ts = slide.shapes.add_table(len(rows)+1, cols, Inches(left), Inches(top), Inches(w), Inches(h))
    table = ts.table
    cw = int(w / cols * 914400)
    for i in range(cols): table.columns[i].width = cw
    for i, h_text in enumerate(headers):
        cell = table.cell(0, i); cell.text = h_text
        for pp in cell.text_frame.paragraphs:
            pp.font.size = Pt(16); pp.font.bold = True; pp.font.color.rgb = TEXT_WHITE
        cell.fill.solid(); cell.fill.fore_color.rgb = ACCENT_BLUE
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r+1, c); cell.text = val
            for pp in cell.text_frame.paragraphs:
                pp.font.size = Pt(14); pp.font.color.rgb = TEXT_LIGHT
            cell.fill.solid(); cell.fill.fore_color.rgb = BG_CARD if r % 2 == 0 else BG_DARK

def content_slide(prs, title, top=1.2):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide); add_bar(slide, 0, 0, 13.333, 0.05, ACCENT_BLUE)
    add_tb(slide, title, left=0.8, top=0.3, w=12, h=0.7, size=28, color=ACCENT_CYAN)
    return slide

def section_cover(prs, num, title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide); add_bar(slide, 0, 0, 13.333, 0.08, ACCENT_CYAN)
    add_tb(slide, f"第{num}部分", left=0.8, top=2.5, w=12, h=0.8, size=24, color=ACCENT_BLUE)
    add_tb(slide, title, left=0.8, top=3.3, w=12, h=1.2, size=40, color=TEXT_WHITE)
    if subtitle: add_body(slide, subtitle, left=0.8, top=4.8, w=12, h=0.6, size=20, color=TEXT_DIM)
    add_bar(slide, 0.8, 6.5, 3, 0.04, ACCENT_CYAN)

# ==================== 模块1：生成式AI概述 ====================
prs1 = new_ppt()

# 封面
slide = prs1.slides.add_slide(prs1.slide_layouts[6])
add_bg(slide); add_bar(slide, 0, 0, 13.333, 0.1, ACCENT_CYAN)
add_tb(slide, "生成式人工智能概述", left=1, top=2.2, w=11, h=1.2, size=44)
add_body(slide, "通识课程 · 模块一", left=1, top=3.6, w=11, h=0.6, size=26, color=ACCENT_BLUE)
add_bar(slide, 1, 4.5, 4, 0.05, ACCENT_CYAN)
add_body(slide, "绍兴柯桥职校骨干教师培训 · 2026年4月", left=1, top=5.2, w=11, h=0.5, size=18, color=TEXT_DIM)

# 什么是生成式AI
slide = content_slide(prs1, "什么是生成式人工智能？")
add_bullets(slide, [
    "📌 定义：生成式AI（Generative AI）是指能够创建新内容（文本、图像、音频、视频等）的人工智能系统",
    "",
    "🔑 核心特征：",
    "  • 从海量数据中学习模式和规律",
    "  • 能够生成原创性内容",
    "  • 支持多模态输入输出",
    "",
    "🔄 与传统AI的区别：",
    "  • 传统AI：分类、预测、识别",
    "  • 生成式AI：创造、生成、设计",
], top=1.2, size=20)

# 发展历程
slide = content_slide(prs1, "发展历程")
add_bullets(slide, [
    "📅 2014年 — GAN（生成对抗网络）诞生",
    "📅 2017年 — Transformer架构提出",
    "📅 2018年 — BERT、GPT-1 发布",
    "📅 2020年 — GPT-3 引发广泛关注",
    "📅 2022年 — DALL·E 2、Stable Diffusion 发布",
    "📅 2022年底 — ChatGPT 引爆全球",
    "📅 2023-2024年 — 多模态大模型爆发",
    "📅 2025年至今 — AI Agent、多模态融合时代",
], top=1.2, size=20)

# 核心技术
slide = content_slide(prs1, "核心技术原理")
add_card(slide, "大语言模型（LLM）",
    "• 基于Transformer架构\n• 自注意力机制\n• 预训练+微调范式", left=0.8, top=1.5, tc=ACCENT_CYAN)
add_card(slide, "扩散模型（Diffusion）",
    "• 图像生成的核心技术\n• 逐步去噪生成过程\n• 从噪声到清晰图像", left=6.8, top=1.5, tc=ACCENT_PURPLE)
add_card(slide, "多模态融合",
    "• 文本、图像、音频统一表示\n• 跨模态理解与生成", left=0.8, top=4.5, tc=ACCENT_GREEN)
add_card(slide, "AI Agent",
    "• 自主规划与执行\n• 工具调用能力\n• 多步骤任务完成", left=6.8, top=4.5, tc=ACCENT_ORANGE)

# 应用场景
slide = content_slide(prs1, "应用场景")
add_bullets(slide, [
    "🎓 教育领域：个性化学习、智能辅导、内容生成",
    "🎨 创意产业：文案写作、设计辅助、视频制作",
    "🏢 企业服务：客服机器人、数据分析、报告生成",
    "🔬 科研创新：文献综述、实验设计、代码生成",
    "🏠 日常生活：智能助手、内容创作、娱乐互动",
], top=1.5, size=22)

# 课堂练习
slide = content_slide(prs1, "🏫 课堂练习")
add_card(slide, "练习1：概念辨析",
    "请用自己的话解释：什么是生成式AI？它与传统的分类/预测AI有什么本质区别？举一个你生活中遇到的生成式AI应用例子。",
    left=0.8, top=1.5, w=11.7, h=1.2, tc=ACCENT_ORANGE)
add_card(slide, "练习2：应用场景 brainstorm",
    "分组讨论：列出你所在学科（语文/数学/英语/专业课等）中，最适合用生成式AI辅助教学的3个场景，并说明理由。",
    left=0.8, top=3.0, w=11.7, h=1.2, tc=ACCENT_ORANGE)
add_card(slide, "练习3：体验对话",
    "打开任意一个AI对话工具（如通义千问），向它提问一个你学科领域的问题，观察回答质量，记录你的感受和发现。",
    left=0.8, top=4.5, w=11.7, h=1.2, tc=ACCENT_ORANGE)

prs1.save('/home/admin/.openclaw/workspace/模块1-生成式AI概述.pptx')
print("✅ 模块1 已保存")

# ==================== 模块2：文本生成基础 ====================
prs2 = new_ppt()

# 封面
slide = prs2.slides.add_slide(prs2.slide_layouts[6])
add_bg(slide); add_bar(slide, 0, 0, 13.333, 0.1, ACCENT_CYAN)
add_tb(slide, "文本生成基础", left=1, top=2.0, w=11, h=1.2, size=44)
add_body(slide, "通识课程 · 模块二 · 提示词设计与优化实战", left=1, top=3.4, w=11, h=0.6, size=24, color=ACCENT_BLUE)
add_bar(slide, 1, 4.3, 4, 0.05, ACCENT_CYAN)
add_body(slide, "绍兴柯桥职校骨干教师培训 · 2026年4月", left=1, top=5.0, w=11, h=0.5, size=18, color=TEXT_DIM)

# 提示词基本概念
slide = content_slide(prs2, "提示词（Prompt）基本概念")
add_bullets(slide, [
    "📌 什么是提示词？",
    "  • 用户与AI模型交互的输入文本",
    "  • 决定了AI生成内容的方向和质量",
    "",
    "🔑 提示词的四大要素：",
    "  • 角色设定：告诉AI它应该扮演什么角色",
    "  • 任务描述：明确要完成的具体任务",
    "  • 约束条件：格式、长度、风格等限制",
    "  • 示例参考：提供期望输出的样例",
], top=1.2, size=20)

# 设计原则
slide = content_slide(prs2, "提示词设计原则")
add_card(slide, "清晰明确", "避免模糊、歧义的表达", left=0.8, top=1.5, tc=ACCENT_CYAN)
add_card(slide, "具体详细", "提供足够的背景信息", left=3.5, top=1.5, tc=ACCENT_BLUE)
add_card(slide, "结构化", "使用清晰的逻辑框架", left=6.2, top=1.5, tc=ACCENT_PURPLE)
add_card(slide, "可迭代", "根据输出结果不断优化", left=8.9, top=1.5, tc=ACCENT_GREEN)
add_card(slide, "示例驱动", "用例子说明期望的输出格式", left=11.6, top=1.5, tc=ACCENT_ORANGE)
add_body(slide, "💡 好的提示词 = 清晰的角色 + 具体的任务 + 明确的约束 + 参考示例",
    left=0.8, top=4.5, size=20, color=ACCENT_CYAN, bold=True)

# 四大技巧
slide = content_slide(prs2, "优化提示词的四大技巧")
add_code(slide, "1. 角色设定法\n你是一位资深语文教师，请为初中生设计一份关于《春》的教案...", left=0.5, top=1.2, w=5.8, h=1.0)
add_code(slide, "2. 任务分解法\n请完成以下三个步骤：\n1. 分析这篇文章的中心思想\n2. 提取3个关键论点\n3. 用简洁的语言总结全文", left=0.5, top=2.6, w=5.8, h=1.5)
add_code(slide, "3. 约束控制法\n请用不超过200字的篇幅，用通俗易懂的语言解释量子计算...", left=6.8, top=1.2, w=6, h=1.0)
add_code(slide, "4. 示例引导法\n参考以下格式，为新产品写一段营销文案：\n示例：【产品名】-【核心卖点】-【使用场景】-【行动号召】", left=6.8, top=2.6, w=6, h=1.5)

# 实例1：教案生成
slide = content_slide(prs2, "实例项目1：课程教案生成")
add_body(slide, "📋 为初中数学\"勾股定理\"生成完整教案", left=0.8, top=1.2, size=18, color=ACCENT_CYAN, bold=True)
add_code(slide, """你是一位有10年教学经验的初中数学教师。请为八年级学生设计一份关于"勾股定理"的教案。

要求：
1. 课时：1课时（45分钟）
2. 包含：教学目标、教学重难点、教学过程、板书设计、作业布置
3. 教学过程要体现"情境导入-探究新知-巩固练习-总结提升"四个环节
4. 至少设计2个学生互动环节
5. 语言简洁专业，适合教师直接使用""", left=0.5, top=2, w=12.3, h=2.8)
add_table(slide, ["维度", "优化前", "优化后"],
    [["角色","无","10年经验的数学教师"],["任务","写教案","具体到课时、环节、互动"],
     ["约束","无","45分钟、4个环节、2个互动"],["输出","模糊","结构化、可直接使用"]], top=5.2, h=1.5)

# 实例2：学生评语
slide = content_slide(prs2, "实例项目2：学生评语生成")
add_body(slide, "📋 为不同特点的学生写个性化期末评语", left=0.8, top=1.2, size=18, color=ACCENT_CYAN, bold=True)
add_code(slide, """你是一位初中班主任。请为以下3位不同特点的学生写期末评语：

学生A：成绩优秀，乐于助人，但体育较弱
学生B：成绩中等，动手能力很强，上课偶尔走神
学生C：成绩落后，但近期进步明显，热爱劳动

要求：
- 每人100-150字
- 先肯定优点，再委婉指出不足
- 语气温暖鼓励，适合家长阅读
- 避免套话，体现个性化""", left=0.5, top=2, w=12.3, h=3.0)

# 实例3：试题生成
slide = content_slide(prs2, "实例项目3：试题生成")
add_body(slide, "📋 为物理章节生成结构化练习题", left=0.8, top=1.2, size=18, color=ACCENT_CYAN, bold=True)
add_code(slide, """你是一位初中物理教师。请为"欧姆定律"章节生成一套练习题：

要求：
- 选择题5道（含答案和解析）
- 填空题3道
- 计算题2道（由易到难）
- 实验探究题1道
- 难度分布：基础60%、中等30%、提高10%
- 标注每道题考查的知识点""", left=0.5, top=2, w=12.3, h=2.5)

# 实例4：教学反思
slide = content_slide(prs2, "实例项目4：教学反思撰写")
add_body(slide, "📋 根据课堂情况撰写结构化教学反思", left=0.8, top=1.2, size=18, color=ACCENT_CYAN, bold=True)
add_code(slide, """你是一位有经验的英语教师。请根据以下课堂情况写一份教学反思：

课堂情况：
- 课题：七年级下册 Unit 5 "What are you doing?"
- 亮点：角色扮演活动学生参与度高
- 不足：语法讲解部分节奏太快，部分学生跟不上
- 改进：下次增加练习环节，使用更多视觉辅助

要求：
- 800字左右
- 结构：教学回顾→成功之处→不足之处→改进措施
- 体现新课标理念""", left=0.5, top=2, w=12.3, h=2.8)

# 实例5：课件大纲
slide = content_slide(prs2, "实例项目5：课件大纲设计")
add_body(slide, "📋 为一节课设计完整的PPT课件大纲", left=0.8, top=1.2, size=18, color=ACCENT_CYAN, bold=True)
add_code(slide, """你是一位教学设计专家。请为初中语文《背影》设计一份PPT课件大纲：

要求：
- 共15-20页PPT
- 包含：导入、作者简介、背景介绍、文本分析、重点段落赏析、
  写作手法、情感教育、课堂讨论、作业布置
- 每页标注：标题、核心内容、配图建议
- 设计3个互动环节（提问/讨论/活动）
- 融入课程思政元素""", left=0.5, top=2, w=12.3, h=2.8)

# 多轮对话
slide = content_slide(prs2, "结合上下文的多轮对话提示词")
add_bullets(slide, [
    "📌 什么是多轮对话提示？",
    "  • 在对话中保持上下文连贯性",
    "  • 基于前一轮的输出进行深化",
    "  • 逐步细化、迭代优化",
    "",
    "🔄 多轮对话示例：",
    "  第一轮：写一段300字的人工智能介绍",
    "  第二轮：重点展开2020-2025年，补充到500字",
    "  第三轮：改写成适合初中生阅读的版本",
    "  第四轮：整理成表格形式（年份/事件/意义）",
], top=1.2, size=18)
add_body(slide, "💡 最佳实践：保持上下文引用 → 渐进式细化 → 及时反馈 → 任务拆分 → 状态保存",
    left=0.8, top=5.5, size=17, color=ACCENT_CYAN, bold=True)

# 结构化提示词框架
slide = content_slide(prs2, "常用结构化提示词框架")
add_card(slide, "CRISPE框架",
    "• Capacity and Role：角色定位\n• Insight：背景信息\n• Statement：任务描述\n• Personality：语言风格\n• Experiment：多角度输出",
    left=0.8, top=1.5, w=5.5, h=3.0, tc=ACCENT_CYAN)
add_card(slide, "CREATE框架",
    "• Context：上下文\n• Request：请求\n• Expectation：期望\n• Action Steps：步骤\n• Type：类型\n• Example：示例",
    left=6.8, top=1.5, w=5.5, h=3.0, tc=ACCENT_PURPLE)
add_card(slide, "模板化提示词",
    "## 任务\n[明确的任务描述]\n## 背景\n[相关背景信息]\n## 要求\n[格式/长度/风格]\n## 示例\n[参考示例]\n## 输出\n[开始生成]",
    left=0.8, top=5.0, w=11.5, h=2.2, tc=ACCENT_GREEN)

# 课堂练习
slide = content_slide(prs2, "🏫 课堂练习")
add_card(slide, "练习1：提示词改写",
    "将以下模糊提示词改写为结构化提示词：\"帮我写一份教案\" → 要求包含角色、任务、约束、示例四大要素。",
    left=0.8, top=1.5, w=11.7, h=1.1, tc=ACCENT_ORANGE)
add_card(slide, "练习2：教案生成实战",
    "选择一个你任教学科的课题，使用\"角色设定+任务分解+约束控制\"技巧，编写一个提示词，让AI生成一份完整教案。",
    left=0.8, top=2.9, w=11.7, h=1.1, tc=ACCENT_ORANGE)
add_card(slide, "练习3：多轮对话优化",
    "用AI工具完成一个多轮对话任务：先让AI写一段文字，然后逐步要求它调整风格、长度、格式，记录每次调整的效果变化。",
    left=0.8, top=4.3, w=11.7, h=1.1, tc=ACCENT_ORANGE)
add_card(slide, "练习4：评语批量生成",
    "为你班级的3位学生写个性化期末评语提示词，要求体现每位学生的特点，避免套话。生成后互相评议质量。",
    left=0.8, top=5.7, w=11.7, h=1.1, tc=ACCENT_ORANGE)

prs2.save('/home/admin/.openclaw/workspace/模块2-文本生成基础.pptx')
print("✅ 模块2 已保存")

# ==================== 模块3：文生图技术 ====================
prs3 = new_ppt()

# 封面
slide = prs3.slides.add_slide(prs3.slide_layouts[6])
add_bg(slide); add_bar(slide, 0, 0, 13.333, 0.1, ACCENT_CYAN)
add_tb(slide, "文生图技术", left=1, top=2.0, w=11, h=1.2, size=44)
add_body(slide, "通识课程 · 模块三 · 从文本到图像的魔法", left=1, top=3.4, w=11, h=0.6, size=24, color=ACCENT_BLUE)
add_bar(slide, 1, 4.3, 4, 0.05, ACCENT_CYAN)
add_body(slide, "绍兴柯桥职校骨干教师培训 · 2026年4月", left=1, top=5.0, w=11, h=0.5, size=18, color=TEXT_DIM)

# 工作原理
slide = content_slide(prs3, "文生图工作原理")
add_bullets(slide, [
    "🔧 核心架构：扩散模型（Diffusion Model）",
    "",
    "📋 工作流程：",
    "  1. 训练阶段：学习图像与文本的对应关系",
    "  2. 生成阶段：从噪声开始，逐步去噪生成图像",
    "  3. 引导过程：文本提示指导生成方向",
    "",
    "⚙️ 关键技术：",
    "  • CLIP模型：文本-图像对齐",
    "  • U-Net架构：图像生成核心",
    "  • 潜在空间扩散：提高效率",
], top=1.2, size=20)

# 常见工具
slide = content_slide(prs3, "常见文生图工具")
add_table(slide, ["工具", "特点", "适用场景"],
    [["Midjourney","艺术性强，质量高","创意设计、艺术创作"],
     ["DALL·E 3","理解力强，安全性好","商业设计、教育素材"],
     ["Stable Diffusion","开源可控，可本地部署","专业创作、定制开发"],
     ["通义万相","中文优化，国内可用","中文内容创作"],
     ["文心一格","中文理解好","中文艺术创作"],
     ["即梦","字节出品，易用","快速设计、社交媒体"]], top=1.5, h=4.5)

# 提示词进阶
slide = content_slide(prs3, "提示词进阶技巧")
add_code(slide, "1. 描述性提示词\n一只穿着宇航服的猫咪，站在月球表面，地球在背景中，\n科幻风格，高清细节，电影级光影", left=0.5, top=1.2, w=6, h=1.3)
add_code(slide, "2. 风格化公式\n[主题]，[艺术风格]，[色彩方案]，[构图方式]，[光影效果]\n\n示例：城市夜景，赛博朋克风格，霓虹灯光，\n俯视构图，高对比度，蓝紫色调", left=0.5, top=2.8, w=6, h=1.8)
add_code(slide, "3. 参数控制\n--ar 16:9          # 宽高比\n--v 6              # 版本\n--style raw        # 原始风格\n--chaos 50         # 创意程度\n--seed 123         # 随机种子", left=6.8, top=1.2, w=6, h=1.8)
add_code(slide, "4. 质量增强词\nmasterpiece, best quality, 8K, ultra detailed,\nprofessional photography, cinematic lighting,\nsharp focus, dramatic atmosphere", left=6.8, top=3.3, w=6, h=1.5)

# 实例1：教育海报
slide = content_slide(prs3, "实例项目1：教育海报设计")
add_body(slide, "📋 为\"校园科技节\"设计一张宣传海报", left=0.8, top=1.2, size=18, color=ACCENT_CYAN, bold=True)
add_code(slide, """设计一张校园科技节宣传海报：
- 主题：探索未来，科技创新
- 风格：现代、科技感、充满活力
- 元素：学生、机器人、太空、未来城市
- 色彩：蓝色和紫色为主，渐变效果
- 构图：中心对称，留白适合添加文字
- 质量：高清，适合打印""", left=0.5, top=2, w=12.3, h=2.5)
add_body(slide, "🔄 迭代过程：初稿→调整科技感元素→优化色彩构图→定稿（添加文字区域留白）", left=0.8, top=4.8, size=17, color=TEXT_DIM)

# 实例2：课文插图
slide = content_slide(prs3, "实例项目2：课文插图生成")
add_body(slide, "📋 为古诗《望庐山瀑布》生成配图", left=0.8, top=1.2, size=18, color=ACCENT_CYAN, bold=True)
add_code(slide, """中国水墨画风格，描绘李白《望庐山瀑布》的意境：
- 高耸的庐山，瀑布从山顶飞流直下
- 远处有阳光照射，彩虹若隐若现
- 近景有松树和岩石
- 留白处适合添加诗句
- 色彩：青绿山水，淡雅清新
- 风格：传统中国画，意境深远""", left=0.5, top=2, w=12.3, h=2.5)

# 实例3：产品概念设计
slide = content_slide(prs3, "实例项目3：产品概念设计")
add_body(slide, "📋 设计一款\"智能学习台灯\"的产品概念图", left=0.8, top=1.2, size=18, color=ACCENT_CYAN, bold=True)
add_code(slide, """产品设计渲染图，智能学习台灯：
- 外观：简约现代，白色主体+木质底座
- 功能元素：LED灯带、触摸屏、摄像头模块
- 场景：放在学生书桌上，旁边有书本和平板
- 光线：温暖柔和的灯光效果
- 风格：苹果风格的产品渲染，干净背景
- 视角：45度角，展示产品全貌
- 质量：8K高清，专业产品摄影""", left=0.5, top=2, w=12.3, h=2.5)

# 实例4：学科知识可视化
slide = content_slide(prs3, "实例项目4：学科知识可视化")
add_body(slide, "📋 生成\"细胞结构\"的科普插图", left=0.8, top=1.2, size=18, color=ACCENT_CYAN, bold=True)
add_code(slide, """生物学教育插图，植物细胞3D结构图：
- 清晰标注：细胞壁、细胞膜、细胞核、叶绿体、线粒体、液泡
- 风格：3D渲染，半透明效果，色彩鲜艳
- 每个细胞器用不同颜色区分
- 背景：浅蓝色渐变
- 适合初中生物课本使用
- 质量：高清，细节丰富""", left=0.5, top=2, w=12.3, h=2.5)

# 实例5：校园风景
slide = content_slide(prs3, "实例项目5：校园风景系列插画")
add_body(slide, "📋 生成\"四季校园\"系列插画", left=0.8, top=1.2, size=18, color=ACCENT_CYAN, bold=True)
add_code(slide, """系列插画：一所中国职业学校的四季景色

春季：樱花树下，教学楼前，学生漫步
夏季：绿树成荫，操场跑道，阳光灿烂
秋季：银杏大道，落叶金黄，图书馆前
冬季：雪景校园，红灯笼，教学楼温暖灯光

统一风格：水彩插画，温暖色调，适合学校宣传册""", left=0.5, top=2, w=12.3, h=2.5)

# 课堂练习
slide = content_slide(prs3, "🏫 课堂练习")
add_card(slide, "练习1：提示词编写",
    "选择一个你学科相关的主题（如古诗词配图、实验示意图、历史场景等），编写一个完整的文生图提示词，包含主题、风格、色彩、构图、质量五个要素。",
    left=0.8, top=1.5, w=11.7, h=1.1, tc=ACCENT_ORANGE)
add_card(slide, "练习2：工具体验",
    "选择一个文生图工具（通义万相/即梦/文心一格），用你编写的提示词生成图片。尝试调整参数（宽高比、风格、创意程度），对比不同参数的效果差异。",
    left=0.8, top=2.9, w=11.7, h=1.1, tc=ACCENT_ORANGE)
add_card(slide, "练习3：迭代优化",
    "对生成的图片进行评估：哪些地方满意？哪些需要改进？修改提示词后重新生成，记录迭代过程和优化思路。",
    left=0.8, top=4.3, w=11.7, h=1.1, tc=ACCENT_ORANGE)
add_card(slide, "练习4：教学应用设计",
    "设计一个将文生图技术融入你课堂教学的方案：在哪个教学环节使用？用来生成什么类型的图片？预期达到什么教学效果？",
    left=0.8, top=5.7, w=11.7, h=1.1, tc=ACCENT_ORANGE)

prs3.save('/home/admin/.openclaw/workspace/模块3-文生图技术.pptx')
print("✅ 模块3 已保存")

# ==================== 模块4：文生视频技术 ====================
prs4 = new_ppt()

# 封面
slide = prs4.slides.add_slide(prs4.slide_layouts[6])
add_bg(slide); add_bar(slide, 0, 0, 13.333, 0.1, ACCENT_CYAN)
add_tb(slide, "文生视频技术", left=1, top=2.0, w=11, h=1.2, size=44)
add_body(slide, "通识课程 · 模块四 · 让文字动起来", left=1, top=3.4, w=11, h=0.6, size=24, color=ACCENT_BLUE)
add_bar(slide, 1, 4.3, 4, 0.05, ACCENT_CYAN)
add_body(slide, "绍兴柯桥职校骨干教师培训 · 2026年4月", left=1, top=5.0, w=11, h=0.5, size=18, color=TEXT_DIM)

# 技术概述
slide = content_slide(prs4, "文生视频技术概述")
add_bullets(slide, [
    "📌 定义：根据文本描述自动生成视频内容",
    "",
    "⚠️ 技术挑战：",
    "  • 时间一致性：帧与帧之间的连贯性",
    "  • 物理规律：符合现实世界的运动规律",
    "  • 语义理解：准确理解复杂文本描述",
    "",
    "🏗️ 主流架构：",
    "  • 扩散模型扩展（Video Diffusion）",
    "  • Transformer架构（Video Transformer）",
    "  • 自回归生成（Autoregressive Generation）",
], top=1.2, size=20)

# 常见工具
slide = content_slide(prs4, "常见文生视频工具")
add_table(slide, ["工具", "特点", "生成时长", "适用场景"],
    [["Sora","高质量，物理规律好","60秒","影视制作、广告"],
     ["Runway Gen-3","控制精细，风格多样","16秒","创意视频、短片"],
     ["Pika","易用，社区活跃","4-7秒","社交媒体、短视频"],
     ["可灵（Kling）","中文优化，国内可用","10秒","中文内容创作"],
     ["智谱清影","国产，性价比高","5秒","教育视频、演示"],
     ["Luma","速度快，质量稳定","5秒","快速原型、测试"]], top=1.5, h=4.0)

# 提示词技巧
slide = content_slide(prs4, "文生视频提示词技巧")
add_code(slide, "1. 镜头语言描述\n[镜头类型] + [运动方式] + [主体] + [环境] + [风格]\n\n示例：特写镜头，缓慢推进，一位科学家在实验室", left=0.5, top=1.2, w=6, h=1.8)
add_code(slide, "2. 时间序列描述\n0-2秒：日出，山间薄雾\n2-4秒：镜头拉远，展现山谷全景\n4-6秒：阳光穿透云层，照亮大地\n6-8秒：鸟群飞过，画面渐暗", left=0.5, top=3.3, w=6, h=1.8)
add_code(slide, "3. 运动控制\n- 摄像机运动：推进、拉远、平移、旋转\n- 主体运动：行走、奔跑、飞行、变形\n- 环境变化：季节更替、天气变化、时间流逝", left=6.8, top=1.2, w=6, h=1.8)
add_code(slide, "4. 风格与质量\n电影级画质，4K分辨率\ncinematic lighting, shallow depth of field\nsmooth motion, natural lighting, realistic", left=6.8, top=3.3, w=6, h=1.5)

# 实例1：科普短视频
slide = content_slide(prs4, "实例项目1：科普短视频制作")
add_body(slide, "📋 制作一段关于\"光合作用\"的科普短视频（15秒）", left=0.8, top=1.2, size=18, color=ACCENT_CYAN, bold=True)
add_code(slide, """制作一段科普短视频，展示植物光合作用的过程：

镜头1（0-5秒）：
微距镜头，阳光照射在绿叶上，叶绿体内部，光反应过程可视化

镜头2（5-10秒）：
动画演示，二氧化碳和水分子进入叶绿体，
葡萄糖和氧气分子生成，化学方程式浮现

镜头3（10-15秒）：
全景镜头，一棵大树在阳光下的生长过程快进，
文字总结：光合作用=光能→化学能

风格：教育动画，清晰易懂，色彩鲜明""", left=0.5, top=2, w=12.3, h=3.0)

# 实例2：历史场景重现
slide = content_slide(prs4, "实例项目2：历史场景重现")
add_body(slide, "📋 制作\"丝绸之路\"历史场景短片", left=0.8, top=1.2, size=18, color=ACCENT_CYAN, bold=True)
add_code(slide, """制作一段丝绸之路历史场景视频（20秒）：

镜头1（0-5秒）：
航拍镜头，敦煌沙漠全景，驼队在沙丘上行进
色调：暖黄色，夕阳时分

镜头2（5-12秒）：
中景镜头，商队到达敦煌古城，城门开启
细节：商人、骆驼、丝绸、香料

镜头3（12-20秒）：
城内集市，各国商人交易，文化交融
风格：电影级画质，历史纪录片风格""", left=0.5, top=2, w=12.3, h=3.0)

# 实例3：物理实验模拟
slide = content_slide(prs4, "实例项目3：物理实验模拟")
add_body(slide, "📋 制作\"电磁感应\"实验模拟视频", left=0.8, top=1.2, size=18, color=ACCENT_CYAN, bold=True)
add_code(slide, """制作电磁感应实验模拟视频（15秒）：

镜头1（0-5秒）：
特写镜头，磁铁插入线圈，电流表指针偏转
风格：3D动画，半透明外壳展示内部结构

镜头2（5-10秒）：
动画演示，磁场线变化，电子流动可视化
色彩：蓝色磁场线，红色电子流

镜头3（10-15秒）：
总结画面，法拉第电磁感应定律公式浮现
风格：教育动画，清晰标注，适合课堂使用""", left=0.5, top=2, w=12.3, h=3.0)

# 实例4：校园宣传片
slide = content_slide(prs4, "实例项目4：校园宣传片")
add_body(slide, "📋 制作职业学校招生宣传片片段", left=0.8, top=1.2, size=18, color=ACCENT_CYAN, bold=True)
add_code(slide, """制作职业学校招生宣传片片段（30秒）：

镜头1（0-8秒）：航拍校园全景，现代化教学楼
镜头2（8-18秒）：实训场景切换（编程/操作设备/实验）
镜头3（18-25秒）：学生活动（社团/竞赛/毕业）
镜头4（25-30秒）：校门口，学生自信走出
文字：欢迎加入，开启你的精彩人生

风格：现代、活力、高质量摄影""", left=0.5, top=2, w=12.3, h=2.5)

# 实例5：古诗词意境视频
slide = content_slide(prs4, "实例项目5：古诗词意境视频")
add_body(slide, "📋 制作《江雪》古诗词意境短片", left=0.8, top=1.2, size=18, color=ACCENT_CYAN, bold=True)
add_code(slide, """制作柳宗元《江雪》意境视频（12秒）：

千山鸟飞绝，万径人踪灭：
航拍镜头，冬日雪山，万籁俱寂
色调：冷蓝色，大雪纷飞

孤舟蓑笠翁，独钓寒江雪：
中景镜头，江面一叶扁舟，老翁垂钓
细节：雪花飘落江面，水面泛起涟漪

风格：中国水墨动画，意境深远
配乐：古琴独奏，空灵悠远""", left=0.5, top=2, w=12.3, h=2.8)

# 课堂练习
slide = content_slide(prs4, "🏫 课堂练习")
add_card(slide, "练习1：镜头语言编写",
    "选择一个你学科相关的知识点（如历史事件、科学实验、文学作品场景等），编写一个分镜头脚本，包含至少3个镜头，每个镜头标注：时长、镜头类型、运动方式、画面内容。",
    left=0.8, top=1.5, w=11.7, h=1.1, tc=ACCENT_ORANGE)
add_card(slide, "练习2：工具体验",
    "选择一个文生视频工具（可灵/智谱清影/Pika），用你编写的分镜头脚本生成视频片段。尝试不同的风格参数，观察效果差异。",
    left=0.8, top=2.9, w=11.7, h=1.1, tc=ACCENT_ORANGE)
add_card(slide, "练习3：教学视频设计",
    "设计一个将文生视频技术融入你课堂教学的方案：在哪个教学环节使用？生成什么类型的视频？如何与现有教学资源结合？",
    left=0.8, top=4.3, w=11.7, h=1.1, tc=ACCENT_ORANGE)
add_card(slide, "练习4：综合创作",
    "小组合作：选择一个主题，分工完成文本脚本→分镜头设计→视频生成→剪辑合成的完整流程，最终展示作品。",
    left=0.8, top=5.7, w=11.7, h=1.1, tc=ACCENT_ORANGE)

prs4.save('/home/admin/.openclaw/workspace/模块4-文生视频技术.pptx')
print("✅ 模块4 已保存")

# ==================== 模块5：跨模型AI拓展能力 ====================
prs5 = new_ppt()

# 封面
slide = prs5.slides.add_slide(prs5.slide_layouts[6])
add_bg(slide); add_bar(slide, 0, 0, 13.333, 0.1, ACCENT_CYAN)
add_tb(slide, "跨模型AI拓展能力", left=1, top=2.0, w=11, h=1.2, size=44)
add_body(slide, "通识课程 · 模块五 · 多模态联动创作实践", left=1, top=3.4, w=11, h=0.6, size=24, color=ACCENT_BLUE)
add_bar(slide, 1, 4.3, 4, 0.05, ACCENT_CYAN)
add_body(slide, "绍兴柯桥职校骨干教师培训 · 2026年4月", left=1, top=5.0, w=11, h=0.5, size=18, color=TEXT_DIM)

# 多模态融合
slide = content_slide(prs5, "多模态融合创作")
add_bullets(slide, [
    "📌 什么是多模态AI？",
    "  • 能够同时处理文本、图像、音频、视频等多种模态",
    "  • 实现跨模态的理解和生成",
    "",
    "📈 技术趋势：",
    "  • 从单模态到多模态",
    "  • 从理解到生成",
    "  • 从独立到融合",
    "",
    "🔄 工作流：",
    "  文本创意 → 图像生成 → 音频配乐 → 视频合成 → 最终作品",
], top=1.2, size=20)

# 联动创作：品牌宣传
slide = content_slide(prs5, "联动创作实例：品牌宣传短片")
add_card(slide, "📝 步骤1：文本创意",
    "• 核心卖点、目标人群、风格定位\n• 使用LLM生成宣传文案", left=0.8, top=1.5, w=5.5, h=1.5, tc=ACCENT_CYAN)
add_card(slide, "🎨 步骤2：图像生成",
    "• 3张场景图：会议室/户外/家庭\n• 使用文生图工具生成", left=6.8, top=1.5, w=5.5, h=1.5, tc=ACCENT_PURPLE)
add_card(slide, "🎵 步骤3：音频生成",
    "• 现代电子+轻音乐，积极向上\n• 使用AI音乐工具生成", left=0.8, top=3.3, w=5.5, h=1.5, tc=ACCENT_GREEN)
add_card(slide, "🎬 步骤4：视频合成",
    "• 整合素材，添加字幕和转场\n• 使用剪辑工具完成", left=6.8, top=3.3, w=5.5, h=1.5, tc=ACCENT_ORANGE)

# 联动创作：教育微课
slide = content_slide(prs5, "联动创作实例：教育微课制作")
add_body(slide, "📋 项目背景：为初中物理课\"电磁感应\"制作一段5分钟微课", left=0.8, top=1.2, size=18, color=ACCENT_CYAN, bold=True)
add_card(slide, "📝 文本脚本",
    "生成微课脚本：引入→原理→实验→应用→总结", left=0.8, top=2.2, w=5.5, h=1.3, tc=ACCENT_CYAN)
add_card(slide, "🎨 视觉素材",
    "法拉第肖像/原理图/实验装置/发电机结构", left=6.8, top=2.2, w=5.5, h=1.3, tc=ACCENT_PURPLE)
add_card(slide, "🎵 配音音频",
    "温和男声，语速适中，轻柔背景音乐", left=0.8, top=3.8, w=5.5, h=1.3, tc=ACCENT_GREEN)
add_card(slide, "🎬 视频合成",
    "PPT+动画+实拍，同步字幕，5分钟", left=6.8, top=3.8, w=5.5, h=1.3, tc=ACCENT_ORANGE)
add_body(slide, "💡 工具链：通义千问（文本）→ 通义万相（图像）→ 剪映（配音）→ 剪映（合成）",
    left=0.8, top=5.5, size=16, color=TEXT_DIM)

# 最佳实践
slide = content_slide(prs5, "跨模型协作最佳实践")
add_card(slide, "1. 统一创意方向",
    "• 先确定核心创意和风格\n• 所有模态围绕同一主题\n• 保持视觉和听觉的一致性", left=0.8, top=1.5, w=3.8, h=2.5, tc=ACCENT_CYAN)
add_card(slide, "2. 迭代优化流程",
    "生成 → 评估 → 调整 → 再生成\n        ↑                    ↓\n        └──── 反馈循环 ────┘", left=5.0, top=1.5, w=3.8, h=2.5, tc=ACCENT_PURPLE)
add_card(slide, "3. 工具链组合",
    "文本：通义千问/GPT-4\n图像：通义万相/Midjourney\n音频：Suno/剪映\n视频：可灵/Runway\n剪辑：剪映/Premiere", left=9.2, top=1.5, w=3.8, h=2.5, tc=ACCENT_GREEN)
add_body(slide, "🔄 核心思路：先定创意方向 → 分模态生成 → 整合优化 → 最终输出",
    left=0.8, top=4.5, size=18, color=ACCENT_CYAN, bold=True)

# 未来展望
slide = content_slide(prs5, "未来展望")
add_card(slide, "技术趋势",
    "• 更高质量的多模态生成\n• 实时交互式创作\n• 个性化定制能力\n• AI Agent自主创作", left=0.8, top=1.5, w=5.5, h=2.5, tc=ACCENT_CYAN)
add_card(slide, "教育应用",
    "• 智能课件自动生成\n• 个性化学习材料\n• 虚拟实验环境\n• 智能评估与反馈", left=6.8, top=1.5, w=5.5, h=2.5, tc=ACCENT_GREEN)
add_card(slide, "挑战与思考",
    "• 版权与伦理问题\n• 内容真实性验证\n• 人机协作边界\n• 教师角色转变", left=0.8, top=4.3, w=11.5, h=2.5, tc=ACCENT_ORANGE)

# 课程总结
slide = content_slide(prs5, "课程总结")
add_bullets(slide, [
    "📌 核心要点回顾：",
    "  1. 生成式AI基础：理解原理，掌握工具",
    "  2. 提示词设计：清晰、具体、结构化",
    "  3. 多模态创作：文本+图像+音频+视频联动",
    "  4. 实践应用：教育、创意、商业场景",
    "",
    "💡 学习建议：",
    "  • 多实践：动手使用各种工具",
    "  • 多迭代：不断优化提示词",
    "  • 多交流：分享经验和技巧",
    "  • 多创新：探索新的应用场景",
], top=1.2, size=20)

# 课堂练习
slide = content_slide(prs5, "🏫 课堂练习")
add_card(slide, "练习1：多模态创作规划",
    "选择一个你学科的教学主题，规划一个多模态创作方案：需要用哪些模态（文本/图像/音频/视频）？每个模态用什么工具？最终产出什么？",
    left=0.8, top=1.5, w=11.7, h=1.1, tc=ACCENT_ORANGE)
add_card(slide, "练习2：工具链实操",
    "选择一个简单主题，完成\"文本→图像→音频\"的完整流程：先用LLM生成文案，再用文生图生成配图，最后用AI工具生成配乐或配音。",
    left=0.8, top=2.9, w=11.7, h=1.1, tc=ACCENT_ORANGE)
add_card(slide, "练习3：教学方案设计",
    "设计一个将多模态AI融入你课堂教学的完整方案，包含：教学目标、教学内容、AI工具选择、教学过程、评价方式。",
    left=0.8, top=4.3, w=11.7, h=1.1, tc=ACCENT_ORANGE)
add_card(slide, "练习4：综合项目展示",
    "小组合作完成一个多模态创作项目（如微课视频、宣传短片等），在班级展示并分享创作过程和心得体会。",
    left=0.8, top=5.7, w=11.7, h=1.1, tc=ACCENT_ORANGE)

# 结束页
slide = prs5.slides.add_slide(prs5.slide_layouts[6])
add_bg(slide); add_bar(slide, 0, 0, 13.333, 0.08, ACCENT_CYAN)
add_tb(slide, "感谢聆听！", left=1, top=2.5, w=11, h=1.2, size=48)
add_bar(slide, 4, 4, 5, 0.05, ACCENT_CYAN)
add_body(slide, "Q&A 环节", left=1, top=4.5, w=11, h=0.8, size=32, color=ACCENT_BLUE)
add_body(slide, "绍兴柯桥职校骨干教师培训 · 生成式人工智能技术应用", left=1, top=5.8, w=11, h=0.5, size=16, color=TEXT_DIM)

prs5.save('/home/admin/.openclaw/workspace/模块5-跨模型AI拓展.pptx')
print("✅ 模块5 已保存")
print("\n🎉 全部5个模块PPT已生成完毕！")
