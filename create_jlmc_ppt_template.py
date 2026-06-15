#!/usr/bin/env python3
"""
吉林医药学院 PPT 模板生成器
生成一套专业的 PPT 模板，包含封面、目录、内容页等
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# 吉林医药学院配色方案
COLORS = {
    'primary': RGBColor(0x00, 0x56, 0x9E),      # 主色调 - 深蓝
    'secondary': RGBColor(0xC4, 0x12, 0x30),     # 辅助色 - 红色
    'accent': RGBColor(0x00, 0x96, 0xD6),        # 强调色 - 亮蓝
    'light_blue': RGBColor(0xE8, 0xF4, 0xFC),    # 浅蓝背景
    'dark': RGBColor(0x33, 0x33, 0x33),          # 深色文字
    'white': RGBColor(0xFF, 0xFF, 0xFF),
    'gray': RGBColor(0x99, 0x99, 0x99),
    'light_gray': RGBColor(0xF5, 0xF5, 0xF5),
}

def create_ppt_template():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    
    # ==================== 第 1 页：封面 ====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
    
    # 背景
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, slide_width, slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = COLORS['white']
    background.line.fill.background()
    
    # 顶部蓝色条
    top_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, slide_width, Inches(2.5)
    )
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = COLORS['primary']
    top_bar.line.fill.background()
    
    # 底部装饰线
    bottom_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(2), Inches(5.5), Inches(9.333), Inches(0.05)
    )
    bottom_line.fill.solid()
    bottom_line.fill.fore_color.rgb = COLORS['secondary']
    bottom_line.line.fill.background()
    
    # 学校名称
    school_name = slide.shapes.add_textbox(
        Inches(2), Inches(0.5), Inches(9.333), Inches(1.5)
    )
    tf = school_name.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = '吉林医药学院'
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER
    
    # 英文校名
    eng_name = slide.shapes.add_textbox(
        Inches(2), Inches(1.8), Inches(9.333), Inches(0.6)
    )
    tf = eng_name.text_frame
    p = tf.paragraphs[0]
    p.text = 'JILIN MEDICAL UNIVERSITY'
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(0xCC, 0xDD, 0xEE)
    p.alignment = PP_ALIGN.CENTER
    
    # 主标题
    title = slide.shapes.add_textbox(
        Inches(1.5), Inches(3.2), Inches(10.333), Inches(1.2)
    )
    tf = title.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = '论文题目/报告主题'
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLORS['dark']
    p.alignment = PP_ALIGN.CENTER
    
    # 副标题/信息
    subtitle = slide.shapes.add_textbox(
        Inches(2), Inches(4.5), Inches(9.333), Inches(1.0)
    )
    tf = subtitle.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = '汇报人：XXX    指导教师：XXX    日期：2026 年 X 月 X 日'
    p.font.size = Pt(18)
    p.font.color.rgb = COLORS['gray']
    p.alignment = PP_ALIGN.CENTER
    
    # 底部装饰条
    bottom_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, Inches(7.2), slide_width, Inches(0.3)
    )
    bottom_bar.fill.solid()
    bottom_bar.fill.fore_color.rgb = COLORS['secondary']
    bottom_bar.line.fill.background()
    
    # ==================== 第 2 页：目录 ====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 背景
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, slide_width, slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['white']
    bg.line.fill.background()
    
    # 顶部装饰条
    top_dec = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, slide_width, Inches(0.08)
    )
    top_dec.fill.solid()
    top_dec.fill.fore_color.rgb = COLORS['primary']
    top_dec.line.fill.background()
    
    # 页面标题
    page_title = slide.shapes.add_textbox(
        Inches(1), Inches(0.5), Inches(11.333), Inches(0.8)
    )
    tf = page_title.text_frame
    p = tf.paragraphs[0]
    p.text = '目录'
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary']
    
    # 标题下划线
    title_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(1), Inches(1.3), Inches(2), Inches(0.04)
    )
    title_line.fill.solid()
    title_line.fill.fore_color.rgb = COLORS['secondary']
    title_line.line.fill.background()
    
    # 目录项
    toc_items = [
        ('01', '研究背景与意义'),
        ('02', '文献综述'),
        ('03', '研究方法与内容'),
        ('04', '研究结果与分析'),
        ('05', '结论与展望'),
    ]
    
    for i, (num, text) in enumerate(toc_items):
        y_pos = Inches(2.0 + i * 0.9)
        
        # 序号
        num_box = slide.shapes.add_textbox(
            Inches(1.5), y_pos, Inches(0.8), Inches(0.6)
        )
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = num
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = COLORS['secondary']
        
        # 文字
        text_box = slide.shapes.add_textbox(
            Inches(2.5), y_pos, Inches(8), Inches(0.6)
        )
        tf = text_box.text_frame
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(24)
        p.font.color.rgb = COLORS['dark']
        
        # 分隔线
        if i < len(toc_items) - 1:
            sep_line = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(1.5), y_pos + Inches(0.65), Inches(9.5), Inches(0.01)
            )
            sep_line.fill.solid()
            sep_line.fill.fore_color.rgb = COLORS['gray']
            sep_line.line.fill.background()
    
    # 页脚
    footer = slide.shapes.add_textbox(
        Inches(11), Inches(7.0), Inches(2), Inches(0.4)
    )
    tf = footer.text_frame
    p = tf.paragraphs[0]
    p.text = '02'
    p.font.size = Pt(12)
    p.font.color.rgb = COLORS['gray']
    p.alignment = PP_ALIGN.RIGHT
    
    # ==================== 第 3 页：内容页模板 ====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 背景
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, slide_width, slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['white']
    bg.line.fill.background()
    
    # 顶部装饰条
    top_dec = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, slide_width, Inches(0.08)
    )
    top_dec.fill.solid()
    top_dec.fill.fore_color.rgb = COLORS['primary']
    top_dec.line.fill.background()
    
    # 章节标题
    section_title = slide.shapes.add_textbox(
        Inches(1), Inches(0.5), Inches(11.333), Inches(0.8)
    )
    tf = section_title.text_frame
    p = tf.paragraphs[0]
    p.text = '研究背景与意义'
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary']
    
    # 标题下划线
    title_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(1), Inches(1.3), Inches(2.5), Inches(0.04)
    )
    title_line.fill.solid()
    title_line.fill.fore_color.rgb = COLORS['secondary']
    title_line.line.fill.background()
    
    # 左侧内容区
    left_content = slide.shapes.add_textbox(
        Inches(1), Inches(1.8), Inches(5.5), Inches(5.0)
    )
    tf = left_content.text_frame
    tf.word_wrap = True
    
    # 添加多个段落
    paragraphs = [
        ('研究背景', True, Pt(22)),
        ('• 随着我国医疗卫生事业的发展，医学教育面临新的挑战和机遇', False, Pt(16)),
        ('• 新医科建设对医学人才培养提出更高要求', False, Pt(16)),
        ('• 本研究立足于吉林医药学院实际，探索创新路径', False, Pt(16)),
        ('', False, Pt(12)),
        ('研究意义', True, Pt(22)),
        ('• 理论意义：丰富医学教育理论体系', False, Pt(16)),
        ('• 实践意义：为院校改革提供参考', False, Pt(16)),
    ]
    
    for i, (text, is_bold, size) in enumerate(paragraphs):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = size
        p.font.bold = is_bold
        p.font.color.rgb = COLORS['dark'] if is_bold else COLORS['dark']
        p.space_after = Pt(6)
    
    # 右侧图表/图片占位区
    placeholder = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.2), Inches(1.8), Inches(5.2), Inches(4.5)
    )
    placeholder.fill.solid()
    placeholder.fill.fore_color.rgb = COLORS['light_blue']
    placeholder.line.color.rgb = COLORS['accent']
    placeholder.line.width = Pt(1.5)
    
    tf = placeholder.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = '图表/图片占位区\n\n建议尺寸：\n宽 5.2 英寸 × 高 4.5 英寸'
    p.font.size = Pt(16)
    p.font.color.rgb = COLORS['primary']
    p.alignment = PP_ALIGN.CENTER
    
    # 页脚
    footer = slide.shapes.add_textbox(
        Inches(11), Inches(7.0), Inches(2), Inches(0.4)
    )
    tf = footer.text_frame
    p = tf.paragraphs[0]
    p.text = '03'
    p.font.size = Pt(12)
    p.font.color.rgb = COLORS['gray']
    p.alignment = PP_ALIGN.RIGHT
    
    # ==================== 第 4 页：两列内容页 ====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 背景
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, slide_width, slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['white']
    bg.line.fill.background()
    
    # 顶部装饰条
    top_dec = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, slide_width, Inches(0.08)
    )
    top_dec.fill.solid()
    top_dec.fill.fore_color.rgb = COLORS['primary']
    top_dec.line.fill.background()
    
    # 页面标题
    page_title = slide.shapes.add_textbox(
        Inches(1), Inches(0.5), Inches(11.333), Inches(0.8)
    )
    tf = page_title.text_frame
    p = tf.paragraphs[0]
    p.text = '文献综述'
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary']
    
    # 标题下划线
    title_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(1), Inches(1.3), Inches(2), Inches(0.04)
    )
    title_line.fill.solid()
    title_line.fill.fore_color.rgb = COLORS['secondary']
    title_line.line.fill.background()
    
    # 左列
    left_col = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(5.5), Inches(5.0)
    )
    left_col.fill.solid()
    left_col.fill.fore_color.rgb = COLORS['light_blue']
    left_col.line.fill.background()
    
    tf = left_col.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = '国内研究现状'
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary']
    
    p = tf.add_paragraph()
    p.text = '\n• 国内学者在医学教育领域的研究成果\n• 主要研究方法和理论框架\n• 存在的问题与不足'
    p.font.size = Pt(16)
    p.font.color.rgb = COLORS['dark']
    p.space_after = Pt(8)
    
    # 右列
    right_col = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.8), Inches(5.5), Inches(5.0)
    )
    right_col.fill.solid()
    right_col.fill.fore_color.rgb = COLORS['light_gray']
    right_col.line.fill.background()
    
    tf = right_col.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = '国外研究现状'
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary']
    
    p = tf.add_paragraph()
    p.text = '\n• 国际医学教育发展趋势\n• 先进教育理念与模式\n• 可借鉴的经验与启示'
    p.font.size = Pt(16)
    p.font.color.rgb = COLORS['dark']
    p.space_after = Pt(8)
    
    # 页脚
    footer = slide.shapes.add_textbox(
        Inches(11), Inches(7.0), Inches(2), Inches(0.4)
    )
    tf = footer.text_frame
    p = tf.paragraphs[0]
    p.text = '04'
    p.font.size = Pt(12)
    p.font.color.rgb = COLORS['gray']
    p.alignment = PP_ALIGN.RIGHT
    
    # ==================== 第 5 页：结束页 ====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 背景
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, slide_width, slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['primary']
    bg.line.fill.background()
    
    # 装饰圆
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(4.5), Inches(1.0), Inches(4.333), Inches(4.333)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = RGBColor(0x00, 0x6B, 0xB8)
    circle.line.fill.background()
    
    # 感谢文字
    thanks = slide.shapes.add_textbox(
        Inches(3), Inches(2.5), Inches(7.333), Inches(1.5)
    )
    tf = thanks.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = '感谢聆听'
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER
    
    # 学校名称
    school = slide.shapes.add_textbox(
        Inches(3), Inches(4.2), Inches(7.333), Inches(0.8)
    )
    tf = school.text_frame
    p = tf.paragraphs[0]
    p.text = '吉林医药学院'
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(0xCC, 0xDD, 0xEE)
    p.alignment = PP_ALIGN.CENTER
    
    # 底部装饰线
    bottom_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(4), Inches(6.0), Inches(5.333), Inches(0.04)
    )
    bottom_line.fill.solid()
    bottom_line.fill.fore_color.rgb = COLORS['secondary']
    bottom_line.line.fill.background()
    
    # 保存
    output_path = '/home/admin/.openclaw/workspace/吉林医药学院_PPT 模板.pptx'
    prs.save(output_path)
    print(f'✅ PPT 模板已生成：{output_path}')
    return output_path

if __name__ == '__main__':
    create_ppt_template()
