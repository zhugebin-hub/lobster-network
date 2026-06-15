#!/usr/bin/env python3
"""Extract text from all PPTX files in a directory and save as JSON."""

import sys
import os
import json
from pptx import Presentation

def extract_ppt_text(ppt_dir, output_json, max_chars=3000):
    """Extract text from all PPTX files in ppt_dir and save to output_json."""
    results = {}
    ppt_files = sorted([f for f in os.listdir(ppt_dir) if f.endswith('.pptx')])
    
    for fname in ppt_files:
        path = os.path.join(ppt_dir, fname)
        try:
            prs = Presentation(path)
            text = []
            slide_count = len(prs.slides)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            t = para.text.strip()
                            if t:
                                text.append(t)
            results[fname] = {
                'slides': slide_count,
                'text': '\n'.join(text)[:max_chars]
            }
            print(f"OK: {fname} ({slide_count} slides)")
        except Exception as e:
            print(f"FAIL: {fname} - {e}")
            results[fname] = {'error': str(e)}
    
    with open(output_json, 'w', encoding='utf-8') as f:
        json.dump(results, f, ensure_ascii=False, indent=2)
    
    print(f"\nTotal PPTs: {len(results)}")
    print(f"Output: {output_json}")

if __name__ == '__main__':
    if len(sys.argv) < 3:
        print("Usage: python3 extract_ppt_text.py <ppt_dir> <output_json>")
        sys.exit(1)
    
    extract_ppt_text(sys.argv[1], sys.argv[2])
