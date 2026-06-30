#!/usr/bin/env python3
"""
Course Teaching Analysis - Main Script
Extracts PPT text, reads student list, and generates analysis data for report generation.
"""

import sys
import os
import json
import openpyxl
from pptx import Presentation

def extract_ppt_texts(ppt_dir, output_json, max_chars=3000):
    """Extract text from all PPTX files."""
    results = {}
    ppt_files = sorted([f for f in os.listdir(ppt_dir) if f.endswith('.pptx')])
    
    for fname in ppt_files:
        path = os.path.join(ppt_dir, fname)
        try:
            prs = Presentation(path)
            text = []
            slide_count = len(prs.slides)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            t = para.text.strip()
                            if t:
                                text.append(t)
            results[fname] = {
                'slides': slide_count,
                'text': '\n'.join(text)[:max_chars]
            }
            print(f"OK: {fname} ({slide_count} slides)")
        except Exception as e:
            print(f"FAIL: {fname} - {e}")
            results[fname] = {'error': str(e)}
    
    with open(output_json, 'w', encoding='utf-8') as f:
        json.dump(results, f, ensure_ascii=False, indent=2)
    
    print(f"\nTotal PPTs: {len(results)}")
    return results

def read_student_list(xlsx_path):
    """Read student list from Excel file."""
    wb = openpyxl.load_workbook(xlsx_path)
    ws = wb[wb.sheetnames[0]]
    
    students = []
    headers = []
    for i, row in enumerate(ws.iter_rows(min_row=1, max_row=ws.max_row, values_only=False)):
        vals = [str(c.value) if c.value is not None else '' for c in row]
        if i == 0:
            headers = vals
        else:
            students.append(dict(zip(headers, vals)))
    
    return students, headers

def analyze_teams(ppt_texts, students):
    """Analyze teams and generate scoring data."""
    # Group students by team
    teams = {}
    for student in students:
        team = student.get('团队名称', '')
        if team and team.strip():
            if team not in teams:
                teams[team] = {'members': [], 'work_name': '', 'captain': ''}
            teams[team]['members'].append(student)
            if student.get('是否队长') == '是':
                teams[team]['captain'] = student.get('姓名', '')
            if student.get('作品名称'):
                teams[team]['work_name'] = student.get('作品名称', '')
    
    # Match teams with PPTs
    team_analysis = []
    for team_name, team_data in teams.items():
        # Find matching PPT
        ppt_match = None
        for ppt_name, ppt_data in ppt_texts.items():
            if any(keyword in ppt_name for keyword in [team_name, team_data['work_name']]):
                ppt_match = ppt_data
                break
        
        team_analysis.append({
            'team_name': team_name,
            'work_name': team_data['work_name'],
            'captain': team_data['captain'],
            'member_count': len(team_data['members']),
            'ppt': ppt_match,
            'has_ppt': ppt_match is not None
        })
    
    return team_analysis

def generate_analysis_data(ppt_dir, xlsx_path, output_json):
    """Generate complete analysis data."""
    print("Extracting PPT texts...")
    ppt_texts = extract_ppt_texts(ppt_dir, output_json.replace('.json', '_ppt.json'))
    
    print("\nReading student list...")
    students, headers = read_student_list(xlsx_path)
    print(f"Total students: {len(students)}")
    
    print("\nAnalyzing teams...")
    team_analysis = analyze_teams(ppt_texts, students)
    
    # Save team analysis
    with open(output_json, 'w', encoding='utf-8') as f:
        json.dump(team_analysis, f, ensure_ascii=False, indent=2)
    
    print(f"\nTotal teams: {len(team_analysis)}")
    print(f"Teams with PPT: {sum(1 for t in team_analysis if t['has_ppt'])}")
    print(f"Output: {output_json}")
    
    return team_analysis

if __name__ == '__main__':
    if len(sys.argv) < 4:
        print("Usage: python3 analyze_course.py <ppt_dir> <xlsx_path> <output.json>")
        sys.exit(1)
    
    generate_analysis_data(sys.argv[1], sys.argv[2], sys.argv[3])
