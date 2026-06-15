#!/usr/bin/env python3
"""Generate PPT presentation from course analysis data."""

import sys
import os
import json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Color scheme
PRIMARY = RGBColor(0x2C, 0x5F, 0xA7)
SECONDARY = RGBColor(0x4A, 0x90, 0xD9)
ACCENT = RGBColor(0xE8, 0x6C, 0x00)
DARK = RGBColor(0x33, 0x33, 0x33)
LIGHT = RGBColor(0xF5, 0xF7, 0xFA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x27, 0xAE, 0x60)
RED = RGBColor(0xE7, 0x4C, 0x3C)

def add_text_box(slide, left, top, width, height, text, font_size=12, font_color=DARK, bold=False, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = font_color
    run.font.bold = bold
    return txBox

def add_table(slide, left, top, width, height, data, font_size=10):
    rows = len(data)
    cols = len(data[0]) if rows > 0 else 0
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    for row_idx in range(rows):
        for col_idx in range(cols):
            cell = table.cell(row_idx, col_idx)
            cell.text = str(data[row_idx][col_idx]) if row_idx < len(data) and col_idx < len(data[0]) else ''
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER
                for run in paragraph.runs:
                    run.font.size = Pt(font_size)
                    run.font.color.rgb = WHITE if row_idx == 0 else DARK
                    run.font.bold = (row_idx == 0)
            if row_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = PRIMARY
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT if row_idx % 2 == 0 else WHITE
    return table

def create_presentation(data, output_path):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = PRIMARY
    
    add_text_box(slide, Inches(1), Inches(2.5), Inches(11.333), Inches(1.5),
        data.get('title', '课程教学分析报告'),
        font_size=36, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    
    add_text_box(slide, Inches(1), Inches(4), Inches(11.333), Inches(1),
        data.get('subtitle', ''),
        font_size=28, font_color=RGBColor(0xBB, 0xD5, 0xED), bold=False, alignment=PP_ALIGN.CENTER)
    
    add_text_box(slide, Inches(1), Inches(5.5), Inches(11.333), Inches(0.8),
        f'📅 {data.get("date", "")}  |  📊 {data.get("summary", "")}',
        font_size=16, font_color=RGBColor(0x99, 0xBB, 0xDD), bold=False, alignment=PP_ALIGN.CENTER)
    
    # Content slides
    for section in data.get('sections', []):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = WHITE
        
        add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
            section['title'], font_size=24, font_color=PRIMARY, bold=True)
        
        y_pos = 1.5
        for item in section.get('items', []):
            if isinstance(item, dict) and 'table' in item:
                add_table(slide, Inches(0.5), Inches(y_pos), Inches(12), Inches(3.5),
                    item['table'], font_size=item.get('font_size', 10))
                y_pos += 4
            else:
                add_text_box(slide, Inches(0.5), Inches(y_pos), Inches(12), Inches(1),
                    str(item), font_size=13, font_color=DARK)
                y_pos += 0.8
    
    # Thank you slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = PRIMARY
    
    add_text_box(slide, Inches(1), Inches(2.5), Inches(11.333), Inches(1.5),
        '感谢聆听！',
        font_size=36, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    
    add_text_box(slide, Inches(1), Inches(4), Inches(11.333), Inches(1),
        f'🦞 虾尔 AI  |  {data.get("date", "")}',
        font_size=20, font_color=RGBColor(0xBB, 0xD5, 0xED), bold=False, alignment=PP_ALIGN.CENTER)
    
    prs.save(output_path)
    print(f"PPT saved: {output_path}")

if __name__ == '__main__':
    if len(sys.argv) < 3:
        print("Usage: python3 generate_ppt.py <data.json> <output.pptx>")
        sys.exit(1)
    
    with open(sys.argv[1], 'r', encoding='utf-8') as f:
        data = json.load(f)
    
    create_presentation(data, sys.argv[2])
