#!/usr/bin/env python3
"""生成毕业论文答辩PPT——人工智能的应用前景与成本及大学生职业规划"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import datetime

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ========== 配色方案 ==========
DARK_BLUE = RGBColor(0x1B, 0x2A, 0x4A)     # 深蓝主色
ACCENT_BLUE = RGBColor(0x2E, 0x86, 0xAB)    # 强调蓝
LIGHT_BLUE = RGBColor(0xEB, 0xF5, 0xFB)      # 浅蓝背景
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
MED_GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
ACCENT_ORANGE = RGBColor(0xE8, 0x6C, 0x00)

# ========== 辅助函数 ==========
def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, fill_color, line_color=None, line_width=None):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.fill.solid()
    else:
        shape.line.fill.background()
    if line_width is not None:
        shape.line.width = Pt(line_width)
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18, bold=False, color=DARK_GRAY, align=PP_ALIGN.LEFT, font_name='微软雅黑'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    return txBox

def add_bullet_slide(slide, title, bullets, subtitle=None, title_color=WHITE, bullet_color=DARK_GRAY, bullet_size=15, spacing=Pt(10)):
    """添加带标题和要点列表的页面"""
    # 标题
    add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11.5), Inches(0.9),
                 title, font_size=28, bold=True, color=title_color)
    # 装饰线
    add_shape(slide, Inches(0.8), Inches(1.15), Inches(2.5), Pt(3), ACCENT_BLUE)
    
    # 副标题
    if subtitle:
        add_text_box(slide, Inches(0.8), Inches(1.3), Inches(11.5), Inches(0.6),
                     subtitle, font_size=14, color=MED_GRAY)
    
    # 要点
    start_y = 1.8 if not subtitle else 2.0
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(start_y), Inches(11.5), Inches(5))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = spacing
        # 添加项目符号
        run = p.add_run(); run.text = '▸  ' + bullet
        run.font.size = Pt(bullet_size)
        run.font.color.rgb = bullet_color
        run.font.name = '微软雅黑'

def add_two_col_slide(slide, title, left_title, left_items, right_title, right_items):
    """两栏布局页面"""
    # 标题
    add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11.5), Inches(0.9),
                 title, font_size=28, bold=True, color=DARK_BLUE)
    add_shape(slide, Inches(0.8), Inches(1.15), Inches(2.5), Pt(3), ACCENT_BLUE)
    
    # 左栏
    add_text_box(slide, Inches(0.8), Inches(1.4), Inches(5.5), Inches(0.5),
                 left_title, font_size=18, bold=True, color=ACCENT_BLUE)
    
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(5.5), Inches(4.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(left_items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(8)
        run = p.add_run(); run.text = '•  ' + item
        run.font.size = Pt(13)
        run.font.color.rgb = DARK_GRAY
        run.font.name = '微软雅黑'
    
    # 右栏
    add_text_box(slide, Inches(7.0), Inches(1.4), Inches(5.5), Inches(0.5),
                 right_title, font_size=18, bold=True, color=ACCENT_BLUE)
    
    txBox = slide.shapes.add_textbox(Inches(7.0), Inches(2.0), Inches(5.5), Inches(4.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(right_items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(8)
        run = p.add_run(); run.text = '•  ' + item
        run.font.size = Pt(13)
        run.font.color.rgb = DARK_GRAY
        run.font.name = '微软雅黑'
    
    # 中间分隔线
    add_shape(slide, Inches(6.6), Inches(1.4), Pt(2), Inches(4.5), LIGHT_GRAY)

def add_section_slide(slide, section_num, section_title):
    """章节过渡页"""
    add_bg(slide, DARK_BLUE)
    # 装饰方块
    add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(7.5), DARK_BLUE)
    # 左侧竖条
    add_shape(slide, Inches(0), Inches(0), Inches(0.15), Inches(7.5), ACCENT_BLUE)
    # 章节号
    add_text_box(slide, Inches(2), Inches(2.2), Inches(4), Inches(1.5),
                 f'PART {section_num}', font_size=20, bold=True, color=ACCENT_BLUE)
    add_text_box(slide, Inches(2), Inches(3.2), Inches(10), Inches(2),
                 section_title, font_size=40, bold=True, color=WHITE)
    # 底部装饰线
    add_shape(slide, Inches(2), Inches(5), Inches(3), Pt(3), ACCENT_BLUE)

# ============================================================
# 幻灯片 1: 封面
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BLUE)
# 左侧竖条
add_shape(slide, Inches(0), Inches(0), Inches(0.2), Inches(7.5), ACCENT_BLUE)

# 学校名称
add_text_box(slide, Inches(1.5), Inches(0.8), Inches(10), Inches(0.8),
             '浙江工商大学', font_size=22, color=RGBColor(0xAA, 0xCC, 0xEE))

# 分割线
add_shape(slide, Inches(1.5), Inches(1.6), Inches(3), Pt(2), ACCENT_BLUE)

# 论文题目主标题
add_text_box(slide, Inches(1.5), Inches(2.2), Inches(10), Inches(1.2),
             '人工智能的应用前景与成本分析', font_size=36, bold=True, color=WHITE)

# 论文题目副标题
add_text_box(slide, Inches(1.5), Inches(3.3), Inches(10), Inches(1),
             '——及其对大学生职业规划的影响研究', font_size=22, color=RGBColor(0xBB, 0xDD, 0xEE))

# 底部信息区
# 信息背景块
add_shape(slide, Inches(1.5), Inches(5), Inches(5), Inches(1.5), RGBColor(0x20, 0x35, 0x55))

info_y = 5.1
add_text_box(slide, Inches(1.7), Inches(info_y), Inches(4.5), Inches(0.4),
             '课程：AI赋能', font_size=14, color=RGBColor(0xCC, 0xDD, 0xEE))
add_text_box(slide, Inches(1.7), Inches(info_y + 0.4), Inches(4.5), Inches(0.4),
             '学院：马克思主义学院  |  专业：宗教学', font_size=14, color=RGBColor(0xCC, 0xDD, 0xEE))
add_text_box(slide, Inches(1.7), Inches(info_y + 0.8), Inches(4.5), Inches(0.4),
             '汇报人：黄友赛', font_size=14, color=RGBColor(0xCC, 0xDD, 0xEE))
add_text_box(slide, Inches(1.7), Inches(info_y + 1.15), Inches(4.5), Inches(0.4),
             '指导教师：诸葛斌  |  2026年6月17日', font_size=14, color=RGBColor(0xCC, 0xDD, 0xEE))

# ============================================================
# 幻灯片 2: 目录
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.8), DARK_BLUE)

add_text_box(slide, Inches(1), Inches(0.4), Inches(10), Inches(1),
             '目  录', font_size=32, bold=True, color=WHITE)
add_shape(slide, Inches(1), Inches(1.3), Inches(2), Pt(3), ACCENT_BLUE)

# 目录项
catalog = [
    ('01', '研究背景与意义', 'AI快速发展带来的机遇与挑战'),
    ('02', 'AI的应用前景', '六大领域的应用现状与趋势'),
    ('03', 'AI的成本分析', '经济、社会、伦理、环境四维'),
    ('04', '对职业规划的影响', '就业市场变革与应对策略'),
    ('05', '高校教育建议', '职业规划教育改革方向'),
    ('06', '结论与展望', '研究总结与未来展望'),
]

for i, (num, title, desc) in enumerate(catalog):
    row = i // 2
    col = i % 2
    x = 1.0 + col * 5.8
    y = 2.2 + row * 1.65
    
    # 序号背景
    add_shape(slide, Inches(x), Inches(y), Inches(0.6), Inches(0.6), ACCENT_BLUE)
    add_text_box(slide, Inches(x), Inches(y + 0.05), Inches(0.6), Inches(0.5),
                 num, font_size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    
    # 标题
    add_text_box(slide, Inches(x + 0.8), Inches(y), Inches(4.5), Inches(0.45),
                 title, font_size=18, bold=True, color=DARK_BLUE)
    
    # 描述
    add_text_box(slide, Inches(x + 0.8), Inches(y + 0.45), Inches(4.5), Inches(0.4),
                 desc, font_size=12, color=MED_GRAY)

# ============================================================
# 幻灯片 3: 研究背景
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.8), DARK_BLUE)
add_text_box(slide, Inches(1), Inches(0.4), Inches(10), Inches(1),
             '研究背景与意义', font_size=32, bold=True, color=WHITE)
add_shape(slide, Inches(1), Inches(1.3), Inches(2), Pt(3), ACCENT_BLUE)

add_two_col_slide(slide, '研究背景与意义',
    '研究背景',
    [
        'AI技术（深度学习、大语言模型、生成式AI）快速发展，从实验室走向大规模产业化',
        '生成式AI每年可为全球经济增加2.6-4.4万亿美元（麦肯锡，2023）',
        'AI技术深刻影响就业市场，引发结构性变革',
    ],
    '研究意义',
    [
        '系统分析AI应用前景与多维度成本，为理解AI影响提供理论框架',
        '探讨AI对就业市场的结构性影响，为大学生职业规划提供参考',
        '提出面向AI时代的职业规划策略，具有重要的现实指导意义',
    ])

# ============================================================
# 幻灯片 4: 章节页 - PART 1
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_slide(slide, '01', 'AI的应用前景')

# ============================================================
# 幻灯片 5: AI应用前景 - 概览
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.8), DARK_BLUE)
add_text_box(slide, Inches(1), Inches(0.4), Inches(10), Inches(1),
             'AI的应用前景 · 六大领域', font_size=32, bold=True, color=WHITE)
add_shape(slide, Inches(1), Inches(1.3), Inches(2), Pt(3), ACCENT_BLUE)

fields = [
    ('医疗健康', '医学影像诊断、药物研发、\n个性化治疗、智能健康监测'),
    ('教育培训', '智能辅导系统、个性化学习、\n教学质量评估'),
    ('金融服务', '风险管理、量化交易、\n智能投顾、AI客服'),
    ('智能制造', '机器视觉检测、预测性维护、\n数字孪生、供应链优化'),
    ('交通运输', '自动驾驶、交通流量优化、\n无人机配送'),
    ('创意内容', '文本生成、图像创作、\n视频生成、内容制作'),
]

for i, (name, desc) in enumerate(fields):
    row = i // 3
    col = i % 3
    x = 0.8 + col * 4.1
    y = 2.2 + row * 2.8
    
    # 卡片背景
    add_shape(slide, Inches(x), Inches(y), Inches(3.7), Inches(2.3), LIGHT_BLUE)
    # 左侧色条
    add_shape(slide, Inches(x), Inches(y), Inches(0.08), Inches(2.3), ACCENT_BLUE)
    # 标题
    add_text_box(slide, Inches(x + 0.3), Inches(y + 0.15), Inches(3.2), Inches(0.5),
                 name, font_size=20, bold=True, color=DARK_BLUE)
    # 描述
    add_text_box(slide, Inches(x + 0.3), Inches(y + 0.7), Inches(3.2), Inches(1.4),
                 desc, font_size=13, color=MED_GRAY)

# ============================================================
# 幻灯片 6: AI应用前景 - 重点数据
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.8), DARK_BLUE)
add_text_box(slide, Inches(1), Inches(0.4), Inches(10), Inches(1),
             'AI应用前景 · 关键数据', font_size=32, bold=True, color=WHITE)
add_shape(slide, Inches(1), Inches(1.3), Inches(2), Pt(3), ACCENT_BLUE)

data_cards = [
    ('2.6-4.4\n万亿美元', '生成式AI年度经济价值增量\n（麦肯锡 2023）'),
    ('36.4%\n年复合增长率', 'AI医疗健康市场规模增速\n（2022-2030, Statista）'),
    ('1.2\n万亿美元', 'AI为金融行业创造价值\n（PwC 2030预测）'),
    ('5000\n亿美元', '全球智能制造市场规模\n（德勤 2025预测）'),
    ('7\n万亿美元', '自动驾驶相关经济产出\n（IHS Markit 2035预测）'),
]

for i, (number, label) in enumerate(data_cards):
    x = 0.8 + i * 2.5
    # 数字
    add_text_box(slide, Inches(x), Inches(2.5), Inches(2.2), Inches(1.5),
                 number, font_size=32, bold=True, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)
    # 说明
    add_text_box(slide, Inches(x), Inches(4.2), Inches(2.2), Inches(1.2),
                 label, font_size=12, color=MED_GRAY, align=PP_ALIGN.CENTER)
    # 分隔线
    if i < len(data_cards) - 1:
        add_shape(slide, Inches(x + 2.4), Inches(3.0), Pt(1), Inches(2), LIGHT_GRAY)

# ============================================================
# 幻灯片 7: 章节页 - PART 2
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_slide(slide, '02', 'AI的成本分析')

# ============================================================
# 幻灯片 8: 经济成本
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.8), DARK_BLUE)
add_text_box(slide, Inches(1), Inches(0.4), Inches(10), Inches(1),
             'AI的成本分析 · 经济成本', font_size=32, bold=True, color=WHITE)
add_shape(slide, Inches(1), Inches(1.3), Inches(2), Pt(3), ACCENT_BLUE)

cost_items = [
    ('研发成本', '训练GPT-4级模型需数千块GPU，成本数千万到上亿美元；\nAI工程师平均年薪超15万美元（美国）'),
    ('部署与运维', 'AI推理需要强大算力支撑，GPT-3每次推理约0.04美元；\n数据存储、清洗和管理成本不容忽视'),
    ('迁移与转型', '超过60%的企业在AI转型中遇到组织文化挑战；\n现有系统改造、流程重设计、员工再培训成本高'),
]

for i, (title, desc) in enumerate(cost_items):
    y = 2.2 + i * 1.7
    # 卡片背景
    add_shape(slide, Inches(0.8), Inches(y), Inches(11.5), Inches(1.4), LIGHT_BLUE)
    # 标题
    add_text_box(slide, Inches(1.2), Inches(y + 0.1), Inches(2), Inches(0.5),
                 title, font_size=18, bold=True, color=DARK_BLUE)
    # 描述
    add_text_box(slide, Inches(3.5), Inches(y + 0.05), Inches(8.5), Inches(1.2),
                 desc, font_size=13, color=MED_GRAY)

# ============================================================
# 幻灯片 9: 社会成本与伦理成本
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.8), DARK_BLUE)
add_text_box(slide, Inches(1), Inches(0.4), Inches(10), Inches(1),
             'AI的成本分析 · 社会与伦理成本', font_size=32, bold=True, color=WHITE)
add_shape(slide, Inches(1), Inches(1.3), Inches(2), Pt(3), ACCENT_BLUE)

add_two_col_slide(slide, '社会成本与伦理成本',
    '社会成本',
    [
        '就业冲击：到2027年，AI将替代8500万岗位，创造9700万新岗位（WEF）',
        '技能错配：被替代与创造的岗位需要完全不同的技能',
        '数字鸿沟：AI专利、投资和人才高度集中于少数国家',
        '再培训成本：全球39%的劳动者需要技能再培训',
    ],
    '伦理与治理成本',
    [
        '隐私保护：大规模数据采集可能导致个人隐私泄露和滥用',
        '算法偏见：AI在招聘、信贷等领域可能存在隐性偏见',
        '责任归属：AI错误决策的责任归属尚无明确法律框架',
        '环境成本：训练大模型的碳排放≈5辆汽车全生命周期',
    ])

# ============================================================
# 幻灯片 10: 章节页 - PART 3
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_slide(slide, '03', '对职业规划的影响')

# ============================================================
# 幻灯片 11: 就业市场结构性变化
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.8), DARK_BLUE)
add_text_box(slide, Inches(1), Inches(0.4), Inches(10), Inches(1),
             '就业市场的结构性变化', font_size=32, bold=True, color=WHITE)
add_shape(slide, Inches(1), Inches(1.3), Inches(2), Pt(3), ACCENT_BLUE)

add_two_col_slide(slide, '就业市场的结构性变化',
    '岗位变化',
    [
        'AI相关专业（CS/数据科学）需求持续增长',
        'AI工程师、数据科学家位居增长最快职业前列',
        '会计、法律助理、翻译等面临自动化替代风险',
        '新职业：AI伦理审查员、提示词工程师等',
    ],
    '技能需求转变',
    [
        '技术技能：编程、数据分析、AI工具使用',
        '软技能：批判性思维、创造力、复杂问题解决',
        '到2027年，分析性思维和创造力最重要（WEF）',
        '"人机协作"成为未来工作新常态',
    ])

# ============================================================
# 幻灯片 12: 大学生面临的挑战
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.8), DARK_BLUE)
add_text_box(slide, Inches(1), Inches(0.4), Inches(10), Inches(1),
             '大学生职业规划面临的挑战', font_size=32, bold=True, color=WHITE)
add_shape(slide, Inches(1), Inches(1.3), Inches(2), Pt(3), ACCENT_BLUE)

challenges = [
    ('01', '职业选择不确定性增加', '热门专业可能在毕业时变冷门，AI迭代速度超出传统职业规划预测能力'),
    ('02', '技能更新持续性要求', '一次学习终身使用的观念已过时，需要培养终身学习的能力和习惯'),
    ('03', '心理压力和焦虑', '学生担心专业被替代，担忧就业前景，需要心理支持和职业辅导'),
]

for i, (num, title, desc) in enumerate(challenges):
    x = 0.8 + i * 4.1
    # 卡片
    add_shape(slide, Inches(x), Inches(2.2), Inches(3.7), Inches(3.5), LIGHT_BLUE)
    add_shape(slide, Inches(x), Inches(2.2), Inches(3.7), Inches(0.08), ACCENT_BLUE)
    # 序号
    add_text_box(slide, Inches(x + 0.2), Inches(2.4), Inches(0.8), Inches(0.5),
                 num, font_size=24, bold=True, color=ACCENT_BLUE)
    # 标题
    add_text_box(slide, Inches(x + 0.2), Inches(2.9), Inches(3.3), Inches(0.5),
                 title, font_size=16, bold=True, color=DARK_BLUE)
    # 描述
    add_text_box(slide, Inches(x + 0.2), Inches(3.5), Inches(3.3), Inches(1.8),
                 desc, font_size=13, color=MED_GRAY)

# ============================================================
# 幻灯片 13: 职业规划策略
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.8), DARK_BLUE)
add_text_box(slide, Inches(1), Inches(0.4), Inches(10), Inches(1),
             '面向AI时代的职业规划策略', font_size=32, bold=True, color=WHITE)
add_shape(slide, Inches(1), Inches(1.3), Inches(2), Pt(3), ACCENT_BLUE)

strategies = [
    '培养跨学科复合能力（"AI+X"模式）',
    '提升数字素养和AI工具使用能力',
    '强化AI难以替代的核心能力（创造力、情感智慧、批判性思维）',
    '制定灵活的职业发展路径（T型/π型知识结构）',
    '重视实践经验和项目经历积累',
]

for i, s in enumerate(strategies):
    y = 2.2 + i * 0.95
    # 序号圆
    add_shape(slide, Inches(1.2), Inches(y + 0.05), Inches(0.5), Inches(0.5), ACCENT_BLUE)
    add_text_box(slide, Inches(1.2), Inches(y + 0.08), Inches(0.5), Inches(0.45),
                 str(i+1), font_size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # 文本
    add_text_box(slide, Inches(2.0), Inches(y + 0.05), Inches(10), Inches(0.6),
                 s, font_size=16, color=DARK_GRAY)

# ============================================================
# 幻灯片 14: AI难以替代的核心能力
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.8), DARK_BLUE)
add_text_box(slide, Inches(1), Inches(0.4), Inches(10), Inches(1),
             'AI难以替代的五大核心能力', font_size=32, bold=True, color=WHITE)
add_shape(slide, Inches(1), Inches(1.3), Inches(2), Pt(3), ACCENT_BLUE)

abilities = [
    ('🎨 创造力', '原创性思维、发散性思维、跨界联想能力'),
    ('💡 情感智慧', '情商、同理心、团队协作、人际沟通'),
    ('🔍 批判性思维', '提出正确问题、复杂情境判断'),
    ('⚖️ 道德判断', '价值决策、伦理考量、责任意识'),
    ('🔄 终身学习', '持续学习、自我更新、适应变化'),
]

for i, (title, desc) in enumerate(abilities):
    x = 0.8 + i * 2.5
    add_shape(slide, Inches(x), Inches(2.2), Inches(2.2), Inches(3.8), LIGHT_BLUE)
    add_shape(slide, Inches(x), Inches(2.2), Inches(2.2), Inches(0.08), ACCENT_BLUE)
    add_text_box(slide, Inches(x + 0.15), Inches(2.6), Inches(1.9), Inches(0.6),
                 title, font_size=16, bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(x + 0.15), Inches(3.4), Inches(1.9), Inches(2),
                 desc, font_size=12, color=MED_GRAY, align=PP_ALIGN.CENTER)

# ============================================================
# 幻灯片 15: 章节页 - PART 4
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_slide(slide, '04', '高校教育建议与结论')

# ============================================================
# 幻灯片 16: 高校教育建议
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.8), DARK_BLUE)
add_text_box(slide, Inches(1), Inches(0.4), Inches(10), Inches(1),
             '高校职业规划教育的建议', font_size=32, bold=True, color=WHITE)
add_shape(slide, Inches(1), Inches(1.3), Inches(2), Pt(3), ACCENT_BLUE)

suggestions = [
    ('AI素养纳入课程', '开设"AI时代的职业发展"等专题课程，提升学生AI认知'),
    ('跨学科培养机制', '设立"AI+X"交叉学科方向，鼓励跨专业选修'),
    ('强化产教融合', '加强校企合作，建立实习基地，引入行业导师'),
    ('个性化辅导服务', '利用AI技术分析学生兴趣和倾向，提供定制化建议'),
    ('关注心理健康', '加强心理健康教育，帮助学生缓解焦虑，保持积极心态'),
]

for i, (title, desc) in enumerate(suggestions):
    y = 2.2 + i * 0.95
    add_shape(slide, Inches(1.2), Inches(y + 0.05), Inches(0.5), Inches(0.5), ACCENT_BLUE)
    add_text_box(slide, Inches(1.2), Inches(y + 0.08), Inches(0.5), Inches(0.45),
                 str(i+1), font_size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(2.0), Inches(y), Inches(3.5), Inches(0.5),
                 title, font_size=16, bold=True, color=DARK_BLUE)
    add_text_box(slide, Inches(5.8), Inches(y + 0.02), Inches(6.5), Inches(0.6),
                 desc, font_size=14, color=MED_GRAY)

# ============================================================
# 幻灯片 17: 结论
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.8), DARK_BLUE)
add_text_box(slide, Inches(1), Inches(0.4), Inches(10), Inches(1),
             '结论', font_size=32, bold=True, color=WHITE)
add_shape(slide, Inches(1), Inches(1.3), Inches(2), Pt(3), ACCENT_BLUE)

conclusions = [
    'AI技术正在深刻改变人类社会，在医疗、教育、金融、制造等领域展现巨大潜力，同时也带来经济、社会、伦理和环境等多维度成本',
    '就业市场正在经历深刻的结构性变革：岗位替代与创造并存，技能需求发生根本性转变，"人机协作"成为未来工作新常态',
    '大学生应主动适应变化，培养跨学科能力、数字素养和AI难以替代的核心能力（创造力、情感智慧、批判性思维等），制定灵活的职业发展路径',
    'AI不是人类的替代者，而是合作伙伴——最重要的能力不是与AI竞争，而是学会与AI协作',
]

add_bullet_slide(slide, '结论', conclusions, bullet_size=15, spacing=Pt(14))

# ============================================================
# 幻灯片 18: 致谢页
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BLUE)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(7.5), DARK_BLUE)
add_shape(slide, Inches(0), Inches(0), Inches(0.2), Inches(7.5), ACCENT_BLUE)

add_text_box(slide, Inches(2), Inches(2.0), Inches(10), Inches(1.5),
             '感谢聆听', font_size=44, bold=True, color=WHITE)

add_shape(slide, Inches(2), Inches(3.3), Inches(3), Pt(3), ACCENT_BLUE)

add_text_box(slide, Inches(2), Inches(3.8), Inches(10), Inches(0.8),
             '敬请各位评委老师批评指正', font_size=22, color=RGBColor(0xAA, 0xCC, 0xEE))

add_text_box(slide, Inches(2), Inches(5.2), Inches(10), Inches(0.5),
             '汇报人：黄友赛  |  指导教师：诸葛斌', font_size=16, color=RGBColor(0x88, 0xAA, 0xCC))
add_text_box(slide, Inches(2), Inches(5.7), Inches(10), Inches(0.5),
             '浙江工商大学 · 马克思主义学院 · 宗教学', font_size=14, color=RGBColor(0x88, 0xAA, 0xCC))

# 底部AI标识
add_text_box(slide, Inches(2), Inches(6.5), Inches(10), Inches(0.4),
             '本文由AI辅助创作', font_size=11, color=RGBColor(0x66, 0x88, 0xAA))

# ========== 保存 ==========
output_path = '/home/admin/.openclaw/workspace/thesis-ai-career/毕业论文答辩PPT_人工智能应用前景与职业规划.pptx'
prs.save(output_path)
print(f'PPT已保存至: {output_path}')
