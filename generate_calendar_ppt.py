#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
简约日历风格 PPT 模板生成器
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# 创建演示文稿
prs = Presentation()
prs.slide_width = Cm(33.867)  # A4 比例 16:9
prs.slide_height = Cm(19.05)

# 定义颜色（简约风格）
COLORS = {
    'primary': (44, 62, 80),      # 深蓝灰
    'secondary': (52, 152, 219),   # 蓝色
    'accent': (231, 76, 60),       # 红色（标记重要日期）
    'light': (236, 240, 241),      # 浅灰背景
    'white': (255, 255, 255),      # 白色
    'text': (44, 62, 80),          # 深灰文字
    'gray': (127, 140, 141),       # 灰色
}

def set_background(slide, color=COLORS['white']):
    """设置幻灯片背景"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*color)

def add_calendar_header(slide, title_text="日历", year="2026", month="4 月"):
    """添加日历风格页眉"""
    # 顶部色带
    left = Cm(0)
    top = Cm(0)
    width = prs.slide_width
    height = Cm(2)
    
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*COLORS['primary'])
    shape.line.fill.background()
    
    # 标题
    textbox = slide.shapes.add_textbox(Cm(1.5), Cm(0.5), Cm(15), Cm(1))
    tf = textbox.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(*COLORS['white'])
    p.font.name = 'Microsoft YaHei'
    
    # 年份月份
    textbox = slide.shapes.add_textbox(Cm(20), Cm(0.6), Cm(12), Cm(0.8))
    tf = textbox.text_frame
    p = tf.paragraphs[0]
    p.text = f"{year}年 {month}"
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(*COLORS['white'])
    p.font.name = 'Microsoft YaHei'
    p.alignment = PP_ALIGN.RIGHT

def add_month_calendar(slide, year=2026, month=4, highlight_day=None):
    """添加月历"""
    # 星期标题
    weekdays = ['一', '二', '三', '四', '五', '六', '日']
    
    # 日历起始位置
    start_x = Cm(2)
    start_y = Cm(5)
    cell_width = Cm(4)
    cell_height = Cm(2.5)
    
    # 计算当月第一天是星期几和总天数
    import calendar
    first_weekday, num_days = calendar.monthrange(year, month)
    # 调整为周一为第一天（Python 中周一=0）
    offset = first_weekday
    
    # 星期标题行
    for i, day in enumerate(weekdays):
        left = start_x + i * (cell_width + Cm(0.3))
        textbox = slide.shapes.add_textbox(left, start_y - Cm(0.8), cell_width, Cm(0.6))
        tf = textbox.text_frame
        p = tf.paragraphs[0]
        p.text = day
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = RGBColor(*COLORS['gray'])
        p.font.name = 'Microsoft YaHei'
        p.alignment = PP_ALIGN.CENTER
    
    # 日期格子
    current_day = 1
    for row in range(6):
        for col in range(7):
            day_num = row * 7 + col - offset
            
            if 1 <= day_num <= num_days:
                left = start_x + col * (cell_width + Cm(0.3))
                top = start_y + row * (cell_height + Cm(0.2))
                
                # 添加格子背景
                if highlight_day and day_num == highlight_day:
                    shape = slide.shapes.add_shape(
                        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, cell_width, cell_height
                    )
                    shape.fill.solid()
                    shape.fill.fore_color.rgb = RGBColor(*COLORS['secondary'])
                    shape.line.fill.background()
                    shape.shadow.inherit = False
                    
                    text_color = COLORS['white']
                else:
                    shape = slide.shapes.add_shape(
                        MSO_SHAPE.RECTANGLE, left, top, cell_width, cell_height
                    )
                    shape.fill.solid()
                    shape.fill.fore_color.rgb = RGBColor(*COLORS['light'])
                    shape.line.fill.background()
                    
                    text_color = COLORS['text']
                
                # 日期数字
                textbox = slide.shapes.add_textbox(left, top, cell_width, cell_height)
                tf = textbox.text_frame
                p = tf.paragraphs[0]
                p.text = str(day_num)
                p.font.size = Pt(20)
                p.font.color.rgb = RGBColor(*text_color)
                p.font.name = 'Microsoft YaHei'
                p.alignment = PP_ALIGN.CENTER
                tf.vertical_anchor = MSO_ANCHOR.MIDDLE

def add_content_placeholder(slide, title="标题", left=Cm(2), top=Cm(8), width=Cm(30), height=Cm(8)):
    """添加内容占位符"""
    # 标题
    textbox = slide.shapes.add_textbox(left, top, width, Cm(1))
    tf = textbox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(*COLORS['text'])
    p.font.name = 'Microsoft YaHei'
    
    # 内容区域
    content_top = top + Cm(1.2)
    content_height = height - Cm(1.2)
    
    textbox = slide.shapes.add_textbox(left, content_top, width, content_height)
    tf = textbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "点击添加内容..."
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(*COLORS['gray'])
    p.font.name = 'Microsoft YaHei'

def add_image_placeholder(slide, left, top, width, height, label="图片占位"):
    """添加图片占位符"""
    # 边框
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*COLORS['light'])
    shape.line.color.rgb = RGBColor(*COLORS['gray'])
    shape.line.dash_style = 4  # 虚线
    shape.line.width = Pt(2)
    
    # 文字提示
    textbox = slide.shapes.add_textbox(left, top + height/2 - Cm(0.5), width, Cm(1))
    tf = textbox.text_frame
    p = tf.paragraphs[0]
    p.text = f"← {label} →"
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(*COLORS['gray'])
    p.font.name = 'Microsoft YaHei'
    p.alignment = PP_ALIGN.CENTER

def add_footer(slide):
    """添加页脚"""
    # 底部线条
    left = Cm(2)
    top = prs.slide_height - Cm(1.5)
    width = prs.slide_width - Cm(4)
    
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, Cm(0.05)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*COLORS['secondary'])
    shape.line.fill.background()
    
    # 页码占位
    textbox = slide.shapes.add_textbox(prs.slide_width - Cm(3), top + Cm(0.3), Cm(2), Cm(0.5))
    tf = textbox.text_frame
    p = tf.paragraphs[0]
    p.text = "页码"
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(*COLORS['gray'])
    p.font.name = 'Microsoft YaHei'
    p.alignment = PP_ALIGN.RIGHT

# ============ 创建幻灯片 ============

# 1. 封面页
slide_layout = prs.slide_layouts[6]  # 空白版式
slide = prs.slides.add_slide(slide_layout)
set_background(slide)

# 封面标题框
title_box = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Cm(0), Cm(0), prs.slide_width, Cm(4)
)
title_box.fill.solid()
title_box.fill.fore_color.rgb = RGBColor(*COLORS['primary'])
title_box.line.fill.background()

# 主标题
textbox = slide.shapes.add_textbox(Cm(2), Cm(1), Cm(20), Cm(2.5))
tf = textbox.text_frame
p = tf.paragraphs[0]
p.text = "工作计划与日程安排"
p.font.size = Pt(40)
p.font.bold = True
p.font.color.rgb = RGBColor(*COLORS['white'])
p.font.name = 'Microsoft YaHei'

# 副标题
textbox = slide.shapes.add_textbox(Cm(2), Cm(3.5), Cm(15), Cm(1))
tf = textbox.text_frame
p = tf.paragraphs[0]
p.text = "简约日历风格演示模板"
p.font.size = Pt(20)
p.font.color.rgb = RGBColor(*COLORS['white'])
p.font.name = 'Microsoft YaHei'

# 装饰性日历元素
add_month_calendar(slide, year=2026, month=4, highlight_day=22)

# 底部信息
textbox = slide.shapes.add_textbox(Cm(2), prs.slide_height - Cm(3), Cm(15), Cm(1))
tf = textbox.text_frame
p = tf.paragraphs[0]
p.text = "汇报人：[您的姓名]"
p.font.size = Pt(14)
p.font.color.rgb = RGBColor(*COLORS['gray'])
p.font.name = 'Microsoft YaHei'

textbox = slide.shapes.add_textbox(Cm(2), prs.slide_height - Cm(2), Cm(15), Cm(1))
tf = textbox.text_frame
p = tf.paragraphs[0]
p.text = "日期：2026 年 4 月"
p.font.size = Pt(14)
p.font.color.rgb = RGBColor(*COLORS['gray'])
p.font.name = 'Microsoft YaHei'

# 2. 目录页
slide = prs.slides.add_slide(slide_layout)
set_background(slide)
add_calendar_header(slide, title_text="目录")

# 目录项
chapters = [
    ("01", "工作计划概览", "本月重点工作任务与目标"),
    ("02", "日程安排", "详细时间节点与里程碑"),
    ("03", "进度跟踪", "完成情况与调整方案"),
    ("04", "总结展望", "成果总结与下月计划"),
]

for i, (num, title, desc) in enumerate(chapters):
    top = Cm(6 + i * 2.5)
    
    # 序号框
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Cm(2), top, Cm(1.5), Cm(1.2)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*COLORS['secondary'])
    shape.line.fill.background()
    
    textbox = slide.shapes.add_textbox(Cm(2), top, Cm(1.5), Cm(1.2))
    tf = textbox.text_frame
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(*COLORS['white'])
    p.font.name = 'Microsoft YaHei'
    p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    # 标题
    textbox = slide.shapes.add_textbox(Cm(4), top, Cm(10), Cm(0.8))
    tf = textbox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(*COLORS['text'])
    p.font.name = 'Microsoft YaHei'
    
    # 描述
    textbox = slide.shapes.add_textbox(Cm(4), top + Cm(0.7), Cm(20), Cm(0.6))
    tf = textbox.text_frame
    p = tf.paragraphs[0]
    p.text = desc
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(*COLORS['gray'])
    p.font.name = 'Microsoft YaHei'

add_footer(slide)

# 3. 月历视图页
slide = prs.slides.add_slide(slide_layout)
set_background(slide)
add_calendar_header(slide, title_text="月度日程", year="2026", month="4 月")

# 添加月历
add_month_calendar(slide, year=2026, month=4, highlight_day=22)

# 右侧备注区
textbox = slide.shapes.add_textbox(Cm(22), Cm(5), Cm(10), Cm(2))
tf = textbox.text_frame
p = tf.paragraphs[0]
p.text = "📌 重要事项"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = RGBColor(*COLORS['accent'])
p.font.name = 'Microsoft YaHei'

textbox = slide.shapes.add_textbox(Cm(22), Cm(6), Cm(10), Cm(8))
tf = textbox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "• 点击添加重要日程\n• 标记关键节点\n• 记录会议安排\n• 跟踪项目进度"
p.font.size = Pt(14)
p.font.color.rgb = RGBColor(*COLORS['text'])
p.font.name = 'Microsoft YaHei'

add_footer(slide)

# 4. 周计划页
slide = prs.slides.add_slide(slide_layout)
set_background(slide)
add_calendar_header(slide, title_text="周工作计划", year="2026", month="第 17 周")

# 周历表格
weekdays_full = ['周一', '周二', '周三', '周四', '周五', '周六', '周日']
start_x = Cm(1.5)
start_y = Cm(4.5)
cell_width = Cm(4.2)
cell_height = Cm(5)

for i, day in enumerate(weekdays_full):
    left = start_x + i * (cell_width + Cm(0.2))
    
    # 星期标题
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, start_y - Cm(0.8), cell_width, Cm(0.6)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*COLORS['primary'])
    shape.line.fill.background()
    
    textbox = slide.shapes.add_textbox(left, start_y - Cm(0.8), cell_width, Cm(0.6))
    tf = textbox.text_frame
    p = tf.paragraphs[0]
    p.text = day
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = RGBColor(*COLORS['white'])
    p.font.name = 'Microsoft YaHei'
    p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    # 内容框
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, start_y, cell_width, cell_height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*COLORS['light'])
    shape.line.color.rgb = RGBColor(*COLORS['gray'])
    shape.line.width = Pt(1)
    
    textbox = slide.shapes.add_textbox(left + Cm(0.2), start_y + Cm(0.2), cell_width - Cm(0.4), cell_height - Cm(0.4))
    tf = textbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "点击添加\n本周工作"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(*COLORS['gray'])
    p.font.name = 'Microsoft YaHei'

add_footer(slide)

# 5. 内容页（带图片）
slide = prs.slides.add_slide(slide_layout)
set_background(slide)
add_calendar_header(slide, title_text="工作内容")

# 左侧文字
add_content_placeholder(slide, title="主要任务", left=Cm(2), top=Cm(5), width=Cm(18), height=Cm(10))

# 右侧图片
add_image_placeholder(slide, left=Cm(21), top=Cm(5), width=Cm(11), height=Cm(8), label="工作照片/图表")

# 底部备注
textbox = slide.shapes.add_textbox(Cm(2), Cm(14), Cm(30), Cm(2))
tf = textbox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "备注：点击文本框可编辑内容，右键图片占位符可插入图片"
p.font.size = Pt(12)
p.font.color.rgb = RGBColor(*COLORS['gray'])
p.font.name = 'Microsoft YaHei'

add_footer(slide)

# 6. 进度跟踪页
slide = prs.slides.add_slide(slide_layout)
set_background(slide)
add_calendar_header(slide, title_text="进度跟踪")

# 进度条示例
tasks = [
    ("任务一：项目启动", 100),
    ("任务二：需求分析", 80),
    ("任务三：方案设计", 60),
    ("任务四：实施执行", 40),
    ("任务五：验收总结", 20),
]

for i, (task, progress) in enumerate(tasks):
    top = Cm(5 + i * 2)
    
    # 任务名称
    textbox = slide.shapes.add_textbox(Cm(2), top, Cm(12), Cm(0.8))
    tf = textbox.text_frame
    p = tf.paragraphs[0]
    p.text = task
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(*COLORS['text'])
    p.font.name = 'Microsoft YaHei'
    
    # 进度条背景
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Cm(15), top, Cm(15), Cm(0.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*COLORS['light'])
    shape.line.fill.background()
    
    # 进度条前景
    if progress > 0:
        progress_width = Cm(15) * progress / 100
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Cm(15), top, progress_width, Cm(0.5)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*COLORS['secondary'])
        shape.line.fill.background()
    
    # 百分比
    textbox = slide.shapes.add_textbox(Cm(31), top, Cm(2), Cm(0.5))
    tf = textbox.text_frame
    p = tf.paragraphs[0]
    p.text = f"{progress}%"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = RGBColor(*COLORS['text'])
    p.font.name = 'Microsoft YaHei'
    p.alignment = PP_ALIGN.RIGHT

add_footer(slide)

# 7. 结束页
slide = prs.slides.add_slide(slide_layout)
set_background(slide)

# 底部色带
shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Cm(0), prs.slide_height - Cm(5), prs.slide_width, Cm(5)
)
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(*COLORS['primary'])
shape.line.fill.background()

# 感谢文字
textbox = slide.shapes.add_textbox(Cm(2), prs.slide_height/2 - Cm(2), Cm(30), Cm(2))
tf = textbox.text_frame
p = tf.paragraphs[0]
p.text = "感谢聆听"
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = RGBColor(*COLORS['text'])
p.font.name = 'Microsoft YaHei'
p.alignment = PP_ALIGN.CENTER

textbox = slide.shapes.add_textbox(Cm(2), prs.slide_height/2 + Cm(0.5), Cm(30), Cm(1))
tf = textbox.text_frame
p = tf.paragraphs[0]
p.text = "敬请批评指正"
p.font.size = Pt(18)
p.font.color.rgb = RGBColor(*COLORS['gray'])
p.font.name = 'Microsoft YaHei'
p.alignment = PP_ALIGN.CENTER

# 装饰日历
add_month_calendar(slide, year=2026, month=4, highlight_day=None)

# 保存文件
output_path = '/home/admin/.openclaw/workspace/简约日历风格 PPT 模板.pptx'
prs.save(output_path)
print(f"PPT 模板已生成：{output_path}")
