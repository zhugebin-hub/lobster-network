#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
生成 1530 安全教育 PPT 课件
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_title_slide(prs, title, subtitle_text):
    """创建标题页"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    
    # 设置标题
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(44)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.name = '微软雅黑'
    
    # 设置副标题
    subtitle_shape = slide.placeholders[1]
    subtitle_shape.text = subtitle_text
    subtitle_shape.text_frame.paragraphs[0].font.size = Pt(24)
    subtitle_shape.text_frame.paragraphs[0].font.name = '微软雅黑'

def create_content_slide(prs, title, content_lines):
    """创建内容页"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    # 设置标题
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(36)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.name = '微软雅黑'
    
    # 设置内容
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()
    
    for i, line in enumerate(content_lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(24)
        p.font.name = '微软雅黑'
        p.space_after = Pt(14)
        if line.startswith('•') or line.startswith('✓') or line.startswith('❌') or line.startswith('✅'):
            p.level = 1

def create_ending_slide(prs, text):
    """创建结束页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白版式
    
    # 添加文本框
    left = Inches(1)
    top = Inches(3)
    width = Inches(8)
    height = Inches(2)
    
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.name = '微软雅黑'

def generate_traffic_safety_ppt():
    """生成交通安全教育 PPT"""
    prs = Presentation()
    
    # 设置幻灯片大小为 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # 第 1 页：封面
    create_title_slide(prs, '🚦 交通安全教育', '珍爱生命 安全出行\n七年级安全教育专题')
    
    # 第 2 页：数据警示
    create_content_slide(prs, '📊 交通事故数据警示', [
        '全国每年中小学生交通事故伤亡：约 20000 人',
        '',
        '主要原因：',
        '• 闯红灯',
        '• 横穿马路',
        '• 骑车违规',
        '• 乘车不安全',
        '',
        '💔 每一个数字背后都是一个家庭！'
    ])
    
    # 第 3 页：步行安全
    create_content_slide(prs, '🚶 步行安全"五要五不要"', [
        '✅ 五要：',
        '• 要走人行道',
        '• 要走斑马线',
        '• 要看信号灯',
        '• 要左右观察',
        '• 要集中注意力',
        '',
        '❌ 五不要：',
        '• 不要闯红灯',
        '• 不要横穿马路',
        '• 不要追逐打闹',
        '• 不要低头看手机',
        '• 不要翻越护栏'
    ])
    
    # 第 4 页：交通标志
    create_content_slide(prs, '🚸 常见交通标志', [
        '⚠️ 警告标志（黄色三角形）：',
        '• 注意行人',
        '• 注意儿童',
        '• 注意信号灯',
        '',
        '🚫 禁令标志（红色圆形）：',
        '• 禁止通行',
        '• 禁止行人',
        '• 禁止鸣喇叭',
        '',
        '➡️ 指示标志（蓝色圆形/方形）：',
        '• 人行横道',
        '• 步行',
        '• 鸣喇叭'
    ])
    
    # 第 5 页：信号灯
    create_content_slide(prs, '🚦 交通信号灯规则', [
        '🔴 红灯：停！禁止通行',
        '',
        '🟢 绿灯：行！可以通行',
        '',
        '🟡 黄灯：等！已越过停止线的可以继续通行',
        '',
        '💡 口诀：',
        '"红灯停，绿灯行，黄灯亮了等一等"',
        '"一停二看三通过"'
    ])
    
    # 第 6 页：乘车安全
    create_content_slide(prs, '🚌 乘车安全须知', [
        '✅ 乘坐公交车：',
        '• 排队候车，先下后上',
        '• 坐稳扶好，头手不伸出窗外',
        '',
        '✅ 乘坐私家车：',
        '• 系好安全带',
        '• 坐后排',
        '• 不干扰驾驶员',
        '',
        '❌ 禁止行为：',
        '• 不坐超载车',
        '• 不坐农用车',
        '• 不坐无牌无证车'
    ])
    
    # 第 7 页：骑行安全
    create_content_slide(prs, '🚴 骑行安全规定', [
        '⚠️ 法律规定：',
        '• 未满 12 周岁不得骑自行车上路',
        '• 未满 16 周岁不得骑电动自行车上路',
        '',
        '✅ 安全骑行：',
        '• 戴好头盔',
        '• 检查车况',
        '• 靠右行驶，不逆行',
        '• 不载人，不追逐',
        '',
        '❌ 危险行为：',
        '• 双手离把',
        '• 戴耳机骑车',
        '• 骑车看手机',
        '• 闯红灯'
    ])
    
    # 第 8 页：校车安全
    create_content_slide(prs, '🚌 校车乘坐规范', [
        '📍 候车时：',
        '• 提前到达，有序排队',
        '• 不追逐打闹',
        '',
        '📍 乘车时：',
        '• 系好安全带，安静坐好',
        '• 听从随车老师管理',
        '',
        '📍 下车时：',
        '• 等车停稳，观察路况',
        '• 从右侧下车'
    ])
    
    # 第 9 页：危险场景
    create_content_slide(prs, '⚠️ 这些行为很危险！', [
        '❌ 在马路上玩滑板/轮滑',
        '❌ 在停车场玩耍',
        '❌ 从停放的车辆中间穿行',
        '❌ 在公交车头前横穿马路',
        '❌ 雨天打伞遮挡视线',
        '❌ 夜间穿深色衣服',
        '',
        '💡 正确做法：',
        '选择安全场所活动，穿着醒目衣物'
    ])
    
    # 第 10 页：恶劣天气
    create_content_slide(prs, '🌧️ 恶劣天气出行安全', [
        '🌧️ 雨天：',
        '• 穿雨衣不打伞',
        '• 慢行防滑',
        '• 避开积水，远离电线杆',
        '',
        '🌫️ 雾天：',
        '• 穿醒目衣服',
        '• 提早出门',
        '• 注意观察，让车辆先过',
        '',
        '❄️ 雪天：',
        '• 穿防滑鞋',
        '• 小步慢行',
        '• 不追逐打闹'
    ])
    
    # 第 11 页：安全承诺
    create_content_slide(prs, '✋ 交通安全承诺', [
        '我承诺：',
        '',
        '✓ 遵守交通法规，文明出行',
        '✓ 过马路走斑马线，不闯红灯',
        '✓ 乘车系安全带，不坐违规车',
        '✓ 未满年龄不骑车，骑车戴头盔',
        '✓ 提醒家人遵守交通规则',
        '',
        '承诺人：___________  日期：___________'
    ])
    
    # 第 12 页：结束页
    create_ending_slide(prs, '🌟 高高兴兴上学 平平安安回家 🌟\n\n交通安全 从我做起！')
    
    # 保存文件
    prs.save('/home/admin/.openclaw/workspace/1530 安全教育/02-课件资源包/PPT/02-交通安全.pptx')
    print('✅ 交通安全.pptx 已生成')

if __name__ == '__main__':
    generate_traffic_safety_ppt()
    print('🎉 PPT 生成完成！')
