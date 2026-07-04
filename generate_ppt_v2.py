#!/usr/bin/env python3
"""生成带图的 PPT: 室外远程驾驶无人车路径规划与避障系统设计
   浅色背景 + 大字 + 嵌入论文图片
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ============ 配色方案（浅色背景） ============
BG_COLOR = RGBColor(0xF5, 0xF7, 0xFA)          # 浅灰白背景
CARD_BG = RGBColor(0xFF, 0xFF, 0xFF)            # 白色卡片
DARK_BLUE = RGBColor(0x1A, 0x3C, 0x6E)          # 深蓝标题
ACCENT_BLUE = RGBColor(0x2E, 0x86, 0xC1)        # 亮蓝
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT = RGBColor(0x2C, 0x3E, 0x50)
GRAY_TEXT = RGBColor(0x7F, 0x8C, 0x8D)
ORANGE = RGBColor(0xE6, 0x7E, 0x22)
GREEN = RGBColor(0x27, 0xAE, 0x60)
LIGHT_BORDER = RGBColor(0xE0, 0xE5, 0xEB)       # 浅边框

IMG_DIR = "/tmp/docx_images"

# ============ 图片映射（基于文档结构和图片尺寸） ============
IMG_MAP = {
    "cover_logo": os.path.join(IMG_DIR, "rId34.jpg"),      # 装饰条
    "fig3_1": os.path.join(IMG_DIR, "rId40.jpg"),          # 智能小车主要模块示意图
    "fig3_2": os.path.join(IMG_DIR, "rId41.jpg"),          # 整体流程图
    "fig3_3": os.path.join(IMG_DIR, "rId46.jpg"),          # 自主循环框图
    "fig4_1": os.path.join(IMG_DIR, "rId47.jpg"),          # Canny预处理流程图
    "fig4_2": os.path.join(IMG_DIR, "rId50.jpg"),          # Canny边缘检测流程
    "fig4_3": os.path.join(IMG_DIR, "rId51.jpg"),          # Canny边缘检测效果
    "fig4_4": os.path.join(IMG_DIR, "rId52.jpg"),          # 变换直线检测结果
    "fig5_1": os.path.join(IMG_DIR, "rId53.jpg"),          # 斑马线识别
    "fig5_2": os.path.join(IMG_DIR, "rId54.jpg"),          # 蓝色挡板识别
    "fig5_3": os.path.join(IMG_DIR, "rId55.jpg"),          # A字母识别
    "fig5_4": os.path.join(IMG_DIR, "rId59.jpg"),          # B字母识别
    "fig5_5": os.path.join(IMG_DIR, "rId61.jpg"),          # 左转向识别
    "fig5_6": os.path.join(IMG_DIR, "rId64.png"),          # 右转向识别
    "fig5_7": os.path.join(IMG_DIR, "rId65.png"),          # PID闭环控制
    "fig5_8": os.path.join(IMG_DIR, "rId66.png"),          # 识别蓝色锥桶
    "fig5_9": os.path.join(IMG_DIR, "rId67.png"),          # 避障程序流程
    "fig2_1": os.path.join(IMG_DIR, "rId37.jpg"),          # 远程操控软件流程图
    "fig2_2": os.path.join(IMG_DIR, "rId38.jpg"),          # 遥控视角图
}

def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_rounded_rect(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def set_text(tf, text, size=18, color=DARK_TEXT, bold=False, alignment=PP_ALIGN.LEFT, font_name="微软雅黑"):
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment

def add_bullet_items(tf, items, size=18, color=DARK_TEXT, spacing=Pt(8), font_name="微软雅黑", indent="  • "):
    tf.clear()
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = indent + item if item else ""
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = spacing

def add_title_bar(slide, title_text, subtitle_text=""):
    """浅色标题栏"""
    bar = add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), DARK_BLUE)
    bar.text_frame.word_wrap = True
    set_text(bar.text_frame, title_text, size=34, color=WHITE, bold=True, alignment=PP_ALIGN.LEFT)
    bar.text_frame.margin_left = Inches(0.7)
    bar.text_frame.margin_top = Inches(0.15)
    
    if subtitle_text:
        sub_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.75), Inches(12), Inches(0.4))
        set_text(sub_box.text_frame, subtitle_text, size=16, color=RGBColor(0xBB,0xCC,0xDD))
    
    # 橙色装饰线
    add_rect(slide, Inches(0), Inches(1.2), Inches(13.333), Inches(0.06), ORANGE)

def add_slide_num(slide, num, total=12):
    txBox = slide.shapes.add_textbox(Inches(12.3), Inches(7.05), Inches(0.8), Inches(0.4))
    set_text(txBox.text_frame, f"{num}/{total}", size=12, color=GRAY_TEXT, alignment=PP_ALIGN.RIGHT)

def add_footer(slide):
    add_rect(slide, Inches(0), Inches(7.3), Inches(13.333), Inches(0.2), DARK_BLUE)

def add_image_with_caption(slide, img_path, left, top, width, height, caption=""):
    """添加图片并可选题注"""
    if not os.path.exists(img_path):
        # 占位框
        box = add_rounded_rect(slide, left, top, width, height, RGBColor(0xEC,0xF0,0xF4), LIGHT_BORDER)
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = "📷 图片缺失"
        p.font.size = Pt(14)
        p.font.color.rgb = GRAY_TEXT
        p.alignment = PP_ALIGN.CENTER
        return
    
    try:
        slide.shapes.add_picture(img_path, left, top, width, height)
    except Exception as e:
        print(f"  ⚠️ 图片插入失败 {img_path}: {e}")
        return
    
    if caption:
        cap_box = slide.shapes.add_textbox(left, top + height + Inches(0.08), width, Inches(0.35))
        set_text(cap_box.text_frame, caption, size=14, color=GRAY_TEXT, alignment=PP_ALIGN.CENTER)


# ============================================
# 第1页：封面
# ============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, RGBColor(0xF0, 0xF4, 0xF8))

# 左侧装饰色块
add_rect(slide, Inches(-0.5), Inches(-0.5), Inches(5.5), Inches(8.5), DARK_BLUE)
add_rect(slide, Inches(0), Inches(0), Inches(5), Inches(0.15), ORANGE)

# 主标题
txBox = slide.shapes.add_textbox(Inches(1.2), Inches(1.5), Inches(4.5), Inches(2.5))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "室外远程驾驶无人车"
p.font.size = Pt(42)
p.font.color.rgb = WHITE
p.font.bold = True
p.font.name = "微软雅黑"

p2 = tf.add_paragraph()
p2.text = "路径规划与避障系统设计"
p2.font.size = Pt(42)
p2.font.color.rgb = ORANGE
p2.font.bold = True
p2.font.name = "微软雅黑"
p2.space_before = Pt(12)

# 副标题
txBox2 = slide.shapes.add_textbox(Inches(1.2), Inches(4.3), Inches(4.5), Inches(0.8))
set_text(txBox2.text_frame, "基于5G通信 · OpenCV · PID控制", size=22, color=RGBColor(0xBB,0xCC,0xDD), font_name="微软雅黑")

# 作者信息
txBox3 = slide.shapes.add_textbox(Inches(1.2), Inches(5.5), Inches(4.5), Inches(1.5))
tf3 = txBox3.text_frame
tf3.word_wrap = True
p3 = tf3.paragraphs[0]
p3.text = "李皓然"
p3.font.size = Pt(22)
p3.font.color.rgb = WHITE
p3.font.name = "微软雅黑"

p4 = tf3.add_paragraph()
p4.text = "电子与电气工程学院"
p4.font.size = Pt(18)
p4.font.color.rgb = RGBColor(0x99,0xAA,0xBB)
p4.font.name = "微软雅黑"
p4.space_before = Pt(8)

p5 = tf3.add_paragraph()
p5.text = "指导教师：孟玲玲"
p5.font.size = Pt(18)
p5.font.color.rgb = RGBColor(0x99,0xAA,0xBB)
p5.font.name = "微软雅黑"
p5.space_before = Pt(4)

p6 = tf3.add_paragraph()
p6.text = "第十九届全国大学生智能汽车竞赛"
p6.font.size = Pt(16)
p6.font.color.rgb = RGBColor(0x77,0x88,0x99)
p6.font.name = "微软雅黑"
p6.space_before = Pt(12)

# 右侧：系统架构图
add_image_with_caption(slide, IMG_MAP["fig3_1"], Inches(6.5), Inches(1.2), Inches(6.2), Inches(4.0), "图1：智能小车主要模块示意图")

# 底部装饰
add_rect(slide, Inches(6.5), Inches(5.6), Inches(6.2), Inches(0.05), LIGHT_BORDER)
txBox4 = slide.shapes.add_textbox(Inches(6.5), Inches(5.8), Inches(6.2), Inches(0.5))
set_text(txBox4.text_frame, "远程遥控 + 自主导航 双模式无人车系统", size=16, color=GRAY_TEXT, alignment=PP_ALIGN.CENTER)


# ============================================
# 第2页：目录
# ============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_COLOR)
add_title_bar(slide, "目  录")

toc_items = [
    ("01", "项目背景与目标"),
    ("02", "系统总体架构"),
    ("03", "5G远程通信技术"),
    ("04", "视觉感知与图像处理"),
    ("05", "标志物识别算法"),
    ("06", "PID控制与避障策略"),
    ("07", "系统总结与展望"),
    ("08", "竞赛历程与致谢"),
]

for i, (num, title) in enumerate(toc_items):
    row = i // 2
    col = i % 2
    x = Inches(1.0 + col * 6.2)
    y = Inches(2.0 + row * 1.3)
    
    card = add_rounded_rect(slide, x, y, Inches(5.6), Inches(1.0), CARD_BG, LIGHT_BORDER)
    
    # 序号
    num_box = slide.shapes.add_textbox(x + Inches(0.3), y + Inches(0.15), Inches(1.0), Inches(0.7))
    set_text(num_box.text_frame, num, size=36, color=ACCENT_BLUE, bold=True, alignment=PP_ALIGN.CENTER)
    
    # 标题
    title_box = slide.shapes.add_textbox(x + Inches(1.5), y + Inches(0.22), Inches(3.8), Inches(0.6))
    set_text(title_box.text_frame, title, size=22, color=DARK_TEXT, bold=True)

add_slide_num(slide, 2)
add_footer(slide)


# ============================================
# 第3页：项目背景与目标
# ============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_COLOR)
add_title_bar(slide, "1. 项目背景与目标")

columns = [
    {
        "title": "面临的挑战",
        "color": RGBColor(0xE7, 0x4C, 0x3C),
        "items": [
            "传统自动驾驶在复杂环境下存在局限",
            "救灾、矿山等危险场景需减少人员风险",
            "信号干扰、数据传输稳定性等难题",
        ]
    },
    {
        "title": "5G带来的机遇",
        "color": ACCENT_BLUE,
        "items": [
            "低延迟：空口时延仅1ms",
            "大带宽：峰值速率达20Gbit/s",
            "高可靠：百万级设备连接/平方公里",
        ]
    },
    {
        "title": "项目目标",
        "color": GREEN,
        "items": [
            "构建高效稳定的5G远程驾驶系统",
            "实现远程遥控 + 自主导航双模式",
            "低成本硬件实现竞赛场景自动驾驶",
        ]
    },
]

for i, col in enumerate(columns):
    x = Inches(0.7 + i * 4.2)
    
    title_card = add_rounded_rect(slide, x, Inches(1.6), Inches(3.8), Inches(0.7), col["color"])
    set_text(title_card.text_frame, col["title"], size=22, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    
    content_card = add_rounded_rect(slide, x, Inches(2.4), Inches(3.8), Inches(3.2), CARD_BG, LIGHT_BORDER)
    tf = content_card.text_frame
    tf.word_wrap = True
    for j, item in enumerate(col["items"]):
        if j == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"  • {item}"
        p.font.size = Pt(18)
        p.font.color.rgb = DARK_TEXT
        p.font.name = "微软雅黑"
        p.space_after = Pt(14)

add_slide_num(slide, 3)
add_footer(slide)


# ============================================
# 第4页：系统总体架构
# ============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_COLOR)
add_title_bar(slide, "2. 系统总体架构", "图3-1 智能小车主要模块示意图")

# 左侧文字
card_left = add_rounded_rect(slide, Inches(0.5), Inches(1.5), Inches(5.5), Inches(5.2), CARD_BG, LIGHT_BORDER)
tf = card_left.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "硬件平台"
p.font.size = Pt(24)
p.font.color.rgb = DARK_BLUE
p.font.bold = True
p.font.name = "微软雅黑"

hardware = [
    "主控：树莓派4B（Raspberry Pi OS）",
    "通信：RM500U 5G模块（免驱即插即用）",
    "感知：USB广角摄像头（320×240，30fps）",
    "执行：数字舵机（68°~92°）+ 直流电机",
    "供电：双独立电池（控制/动力隔离）",
]
for item in hardware:
    p = tf.add_paragraph()
    p.text = f"  • {item}"
    p.font.size = Pt(18)
    p.font.color.rgb = DARK_TEXT
    p.font.name = "微软雅黑"
    p.space_after = Pt(10)

p = tf.add_paragraph()
p.text = ""
p.font.size = Pt(6)

p = tf.add_paragraph()
p.text = "软件架构（Python模块化）"
p.font.size = Pt(24)
p.font.color.rgb = DARK_BLUE
p.font.bold = True
p.font.name = "微软雅黑"
p.space_before = Pt(8)

software = [
    "图像采集 → 图像处理 → 标志物识别",
    "PID控制 → 通信 → 语音播报",
]
for item in software:
    p = tf.add_paragraph()
    p.text = f"  • {item}"
    p.font.size = Pt(18)
    p.font.color.rgb = DARK_TEXT
    p.font.name = "微软雅黑"
    p.space_after = Pt(10)

# 右侧图片
add_image_with_caption(slide, IMG_MAP["fig3_1"], Inches(6.3), Inches(1.5), Inches(6.5), Inches(3.8), "图3-1：智能小车主要模块示意图")

add_slide_num(slide, 4)
add_footer(slide)


# ============================================
# 第5页：5G远程通信
# ============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_COLOR)
add_title_bar(slide, "3. 5G远程通信与视觉技术基础", "图2-1 远程操控软件流程图 | 图2-2 遥控视角图")

# 左侧文字
card = add_rounded_rect(slide, Inches(0.5), Inches(1.5), Inches(5.5), Inches(5.2), CARD_BG, LIGHT_BORDER)
tf = card.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "5G通信模块"
p.font.size = Pt(24)
p.font.color.rgb = DARK_BLUE
p.font.bold = True
p.font.name = "微软雅黑"

items_5g = [
    "车载视频回传 + 控制指令下发",
    "JT/T808协议封装 → 5G发送",
    "网页控制界面实时显示",
    "RM500U模块免驱即插即用",
]
for item in items_5g:
    p = tf.add_paragraph()
    p.text = f"  • {item}"
    p.font.size = Pt(18)
    p.font.color.rgb = DARK_TEXT
    p.font.name = "微软雅黑"
    p.space_after = Pt(10)

p = tf.add_paragraph()
p.text = ""
p.font.size = Pt(6)

p = tf.add_paragraph()
p.text = "OpenCV视觉处理"
p.font.size = Pt(24)
p.font.color.rgb = DARK_BLUE
p.font.bold = True
p.font.name = "微软雅黑"
p.space_before = Pt(8)

items_cv = [
    "开源跨平台计算机视觉库",
    "树莓派嵌入式实时视觉任务",
    "采集→灰度→二值化→滤波→边缘→霍夫",
]
for item in items_cv:
    p = tf.add_paragraph()
    p.text = f"  • {item}"
    p.font.size = Pt(18)
    p.font.color.rgb = DARK_TEXT
    p.font.name = "微软雅黑"
    p.space_after = Pt(10)

# 右侧两张图
add_image_with_caption(slide, IMG_MAP["fig2_1"], Inches(6.3), Inches(1.5), Inches(6.5), Inches(2.8), "图2-1：远程操控软件流程图")
add_image_with_caption(slide, IMG_MAP["fig2_2"], Inches(6.3), Inches(4.5), Inches(6.5), Inches(2.2), "图2-2：遥控视角图")

add_slide_num(slide, 5)
add_footer(slide)


# ============================================
# 第6页：图像预处理流水线
# ============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_COLOR)
add_title_bar(slide, "4. 视觉感知 – 图像预处理流水线", "图4-1 Canny边缘检测预处理流程 | 图4-3 Canny边缘检测效果")

# 流程步骤（6步）
steps = [
    ("① ROI裁剪", "保留下方3/4\n排除干扰"),
    ("② 灰度化", "BGR→单通道\n数据量降至1/3"),
    ("③ 二值化", "固定+自适应\n8×6子块"),
    ("④ 降噪", "双边滤波+高斯\n保边去噪"),
    ("⑤ Canny", "滞后阈值30/100\n提取锐利边缘"),
    ("⑥ 形态学", "膨胀→开→闭\n连接断裂"),
]

for i, (title, desc) in enumerate(steps):
    x = Inches(0.5 + i * 2.0)
    y = Inches(1.6)
    
    card = add_rounded_rect(slide, x, y, Inches(1.8), Inches(1.3), ACCENT_BLUE)
    tf = card.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(18)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = "微软雅黑"
    p.alignment = PP_ALIGN.CENTER
    
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(13)
    p2.font.color.rgb = RGBColor(0xDD,0xEE,0xFF)
    p2.font.name = "微软雅黑"
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(4)
    
    if i < len(steps) - 1:
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x + Inches(1.85), y + Inches(0.4), Inches(0.25), Inches(0.4))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = ORANGE
        arrow.line.fill.background()

# 底部图片
add_image_with_caption(slide, IMG_MAP["fig4_1"], Inches(0.5), Inches(3.2), Inches(6.0), Inches(3.0), "图4-1：Canny边缘检测预处理流程")
add_image_with_caption(slide, IMG_MAP["fig4_3"], Inches(6.8), Inches(3.2), Inches(6.0), Inches(3.0), "图4-3：Canny边缘检测效果对比")

add_slide_num(slide, 6)
add_footer(slide)


# ============================================
# 第7页：车道线检测与中线提取
# ============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_COLOR)
add_title_bar(slide, "4. 视觉感知 – 车道线检测与中线提取", "图4-4 变换直线检测结果")

# 三列
cols = [
    ("霍夫变换参数", [
        "rho=1, theta=1°",
        "threshold=70",
        "minLineLength=25",
        "maxLineGap=5",
        "",
        "角度过滤：",
        "仅保留10°~170°",
        "去除水平/垂直线",
    ]),
    ("中线提取算法", [
        "形态学膨胀补间隙",
        "自中间列向左右扫描",
        "检测白像素跳变",
        "取左右边界平均值",
        "",
        "横向偏差 = PID输入",
    ]),
    ("稳定性保障", [
        "历史记忆补偿",
        "中值滤波（5帧）",
        "滑动平均（窗口3）",
        "低通滤波+死区",
    ]),
]

for i, (title, items) in enumerate(cols):
    x = Inches(0.5 + i * 4.2)
    card = add_rounded_rect(slide, x, Inches(1.5), Inches(3.8), Inches(3.0), CARD_BG, LIGHT_BORDER)
    tf = card.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(22)
    p.font.color.rgb = DARK_BLUE
    p.font.bold = True
    p.font.name = "微软雅黑"
    
    for item in items:
        p = tf.add_paragraph()
        p.text = f"  • {item}" if item else ""
        p.font.size = Pt(17)
        p.font.color.rgb = DARK_TEXT
        p.font.name = "微软雅黑"
        p.space_after = Pt(6)

# 底部图片
add_image_with_caption(slide, IMG_MAP["fig4_4"], Inches(0.5), Inches(4.7), Inches(12.3), Inches(2.3), "图4-4：霍夫直线检测与赛道中线提取结果")

add_slide_num(slide, 7)
add_footer(slide)


# ============================================
# 第8页：标志物识别（一）斑马线+蓝色挡板
# ============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_COLOR)
add_title_bar(slide, "5. 标志物识别算法（一）", "图5-1 斑马线识别 | 图5-2 蓝色挡板识别")

# 左侧文字
card_left = add_rounded_rect(slide, Inches(0.5), Inches(1.5), Inches(5.5), Inches(2.5), CARD_BG, LIGHT_BORDER)
tf = card_left.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "斑马线识别"
p.font.size = Pt(24)
p.font.color.rgb = DARK_BLUE
p.font.bold = True
p.font.name = "微软雅黑"

for item in ["高阈值提取白色区域 → 轮廓查找", "筛选：面积合理 + 长宽比2:1~5:1 + 数量≥3", "触发后：停车5秒 + 语音播报", "准确率 > 98%"]:
    p = tf.add_paragraph()
    p.text = f"  • {item}"
    p.font.size = Pt(17)
    p.font.color.rgb = DARK_TEXT
    p.font.name = "微软雅黑"
    p.space_after = Pt(6)

# 右侧文字
card_right = add_rounded_rect(slide, Inches(6.3), Inches(1.5), Inches(6.5), Inches(2.5), CARD_BG, LIGHT_BORDER)
tf2 = card_right.text_frame
tf2.word_wrap = True
p = tf2.paragraphs[0]
p.text = "蓝色挡板识别（发车/停车信号）"
p.font.size = Pt(24)
p.font.color.rgb = DARK_BLUE
p.font.bold = True
p.font.name = "微软雅黑"

for item in ["HSV蓝色范围[110,50,50]~[130,255,255]", "inRange生成蓝色掩膜 → 腐蚀/膨胀", "面积 > 图像10% → 判定挡板", "挡板消失连续3帧 → 自动发车"]:
    p = tf2.add_paragraph()
    p.text = f"  • {item}"
    p.font.size = Pt(17)
    p.font.color.rgb = DARK_TEXT
    p.font.name = "微软雅黑"
    p.space_after = Pt(6)

# 底部两张图
add_image_with_caption(slide, IMG_MAP["fig5_1"], Inches(0.5), Inches(4.2), Inches(6.0), Inches(2.8), "图5-1：斑马线识别")
add_image_with_caption(slide, IMG_MAP["fig5_2"], Inches(6.8), Inches(4.2), Inches(6.0), Inches(2.8), "图5-2：蓝色挡板识别")

add_slide_num(slide, 8)
add_footer(slide)


# ============================================
# 第9页：标志物识别（二）字母+变道
# ============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_COLOR)
add_title_bar(slide, "5. 标志物识别算法（二）", "图5-3 A字母识别 | 图5-4 B字母识别 | 图5-5 左转向 | 图5-6 右转向")

# 四张图
add_image_with_caption(slide, IMG_MAP["fig5_3"], Inches(0.5), Inches(1.5), Inches(6.0), Inches(2.5), "图5-3：A字母识别（顶点数=3，形状因子<0.25）")
add_image_with_caption(slide, IMG_MAP["fig5_4"], Inches(6.8), Inches(1.5), Inches(6.0), Inches(2.5), "图5-4：B字母识别（顶点数=4，形状因子1.5~2.1）")
add_image_with_caption(slide, IMG_MAP["fig5_5"], Inches(0.5), Inches(4.3), Inches(6.0), Inches(2.5), "图5-5：左转向识别")
add_image_with_caption(slide, IMG_MAP["fig5_6"], Inches(6.8), Inches(4.3), Inches(6.0), Inches(2.5), "图5-6：右转向识别")

add_slide_num(slide, 9)
add_footer(slide)


# ============================================
# 第10页：PID控制与避障策略
# ============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_COLOR)
add_title_bar(slide, "6. PID控制与避障策略", "图5-7 PID闭环控制 | 图5-8 蓝色锥桶 | 图5-9 避障流程")

# 三列文字
cols = [
    ("舵机转向 — PD控制", [
        "输出 = Kp×偏差 + Kd×Δ偏差",
        "去除积分项防振荡",
        "PWM 50Hz，角度68°~92°",
    ]),
    ("电机速度 — 增量式PID", [
        "增量限制±10",
        "占空比30%~90%",
        "防速度突变打滑",
    ]),
    ("避障策略", [
        "障碍物>30% → 减速60%",
        "正前方>50% → 后退0.8s",
        "视觉→舵机延迟 < 80ms",
    ]),
]

for i, (title, items) in enumerate(cols):
    x = Inches(0.5 + i * 4.2)
    card = add_rounded_rect(slide, x, Inches(1.5), Inches(3.8), Inches(2.5), CARD_BG, LIGHT_BORDER)
    tf = card.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(20)
    p.font.color.rgb = DARK_BLUE
    p.font.bold = True
    p.font.name = "微软雅黑"
    
    for item in items:
        p = tf.add_paragraph()
        p.text = f"  • {item}"
        p.font.size = Pt(17)
        p.font.color.rgb = DARK_TEXT
        p.font.name = "微软雅黑"
        p.space_after = Pt(8)

# 底部三张图
add_image_with_caption(slide, IMG_MAP["fig5_7"], Inches(0.5), Inches(4.3), Inches(3.8), Inches(2.5), "图5-7：PID闭环控制示意图")
add_image_with_caption(slide, IMG_MAP["fig5_8"], Inches(4.75), Inches(4.3), Inches(3.8), Inches(2.5), "图5-8：识别蓝色锥桶")
add_image_with_caption(slide, IMG_MAP["fig5_9"], Inches(9.0), Inches(4.3), Inches(3.8), Inches(2.5), "图5-9：避障程序流程")

add_slide_num(slide, 10)
add_footer(slide)


# ============================================
# 第11页：系统总结与展望
# ============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_COLOR)
add_title_bar(slide, "7. 系统总结与展望")

# 左侧：成果
card1 = add_rounded_rect(slide, Inches(0.5), Inches(1.5), Inches(7.5), Inches(5.0), CARD_BG, LIGHT_BORDER)
tf = card1.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "已完成成果"
p.font.size = Pt(26)
p.font.color.rgb = GREEN
p.font.bold = True
p.font.name = "微软雅黑"

results = [
    "✅ 树莓派4B + 5G + OpenCV无人车系统",
    "✅ 远程遥控模式：5G视频回传 + 网页操控",
    "✅ 自主导航模式：视觉循迹 + 标志物识别 + 避障",
    "✅ 完整视觉流水线：ROI→灰度→二值化→Canny→霍夫→中线",
    "✅ 标志物识别：斑马线/挡板/A字母/B字母/变道线",
    "✅ PID控制 + 有限状态机（5状态自动切换）",
    "✅ 光照鲁棒性：自适应阈值 + 20+组参数组合",
]
for item in results:
    p = tf.add_paragraph()
    p.text = f"  {item}"
    p.font.size = Pt(19)
    p.font.color.rgb = DARK_TEXT
    p.font.name = "微软雅黑"
    p.space_after = Pt(10)

# 右侧：展望
card2 = add_rounded_rect(slide, Inches(8.3), Inches(1.5), Inches(4.5), Inches(5.0), DARK_BLUE)
tf2 = card2.text_frame
tf2.word_wrap = True
p = tf2.paragraphs[0]
p.text = "未来展望"
p.font.size = Pt(26)
p.font.color.rgb = ORANGE
p.font.bold = True
p.font.name = "微软雅黑"

future = [
    ("多传感器融合", "激光雷达 + IMU + GPS"),
    ("深度学习", "目标检测替代传统视觉"),
    ("应用场景", "巡检 · 物流 · 农业"),
    ("性能优化", "更高帧率 · 更低延迟"),
]
for title, desc in future:
    p = tf2.add_paragraph()
    p.text = f"  {title}"
    p.font.size = Pt(20)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = "微软雅黑"
    p.space_before = Pt(14)
    
    p2 = tf2.add_paragraph()
    p2.text = f"  {desc}"
    p2.font.size = Pt(17)
    p2.font.color.rgb = RGBColor(0xBB,0xCC,0xDD)
    p2.font.name = "微软雅黑"

add_slide_num(slide, 11)
add_footer(slide)


# ============================================
# 第12页：竞赛历程与致谢
# ============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_COLOR)
add_title_bar(slide, "8. 竞赛历程与致谢")

# 左侧：历程
card1 = add_rounded_rect(slide, Inches(0.5), Inches(1.5), Inches(6.0), Inches(5.0), CARD_BG, LIGHT_BORDER)
tf = card1.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "备赛历程（5个月）"
p.font.size = Pt(26)
p.font.color.rgb = DARK_BLUE
p.font.bold = True
p.font.name = "微软雅黑"

journey = [
    "6人团队：机械/硬件/通信/视觉/控制/网页",
    "从零搭建：刷系统→GPIO→摄像头→5G配置",
    "参数调优：记录20+组天气/时段参数组合",
    "深夜抢修：舵机齿轮崩裂，凌晨3点重新上路",
]
for item in journey:
    p = tf.add_paragraph()
    p.text = f"  • {item}"
    p.font.size = Pt(19)
    p.font.color.rgb = DARK_TEXT
    p.font.name = "微软雅黑"
    p.space_after = Pt(14)

p = tf.add_paragraph()
p.text = ""
p.font.size = Pt(8)

p = tf.add_paragraph()
p.text = "心得体会"
p.font.size = Pt(24)
p.font.color.rgb = DARK_BLUE
p.font.bold = True
p.font.name = "微软雅黑"
p.space_before = Pt(12)

lessons = [
    "书本公式与工程应用差距巨大",
    "良好分工和信任是复杂项目的基础",
]
for item in lessons:
    p = tf.add_paragraph()
    p.text = f"  • {item}"
    p.font.size = Pt(18)
    p.font.color.rgb = DARK_TEXT
    p.font.name = "微软雅黑"
    p.space_after = Pt(10)

# 右侧：致谢
card2 = add_rounded_rect(slide, Inches(6.8), Inches(1.5), Inches(6.0), Inches(5.0), DARK_BLUE)
tf2 = card2.text_frame
tf2.word_wrap = True
p = tf2.paragraphs[0]
p.text = "致谢"
p.font.size = Pt(28)
p.font.color.rgb = ORANGE
p.font.bold = True
p.font.name = "微软雅黑"

thanks = [
    "感谢导师孟玲玲老师的悉心指导",
    "",
    "感谢队友胡应浩、崔浩宇、",
    "韩磊、刘泽悦的并肩作战！",
    "",
    "五个月备赛，从陌生到默契，",
    "从失败到成功。",
    "",
    "这段经历，一生珍藏。",
]
for item in thanks:
    p = tf2.add_paragraph()
    p.text = f"  {item}" if item else ""
    p.font.size = Pt(19)
    p.font.color.rgb = WHITE
    p.font.name = "微软雅黑"
    p.space_after = Pt(4)

add_slide_num(slide, 12)
add_footer(slide)


# ============================================
# 保存
# ============================================
output_path = "/home/admin/.openclaw/workspace/室外远程驾驶无人车路径规划与避障系统设计_李皓然_带图版.pptx"
prs.save(output_path)
print(f"PPT已保存: {output_path}")
print(f"文件大小: {os.path.getsize(output_path) / 1024:.1f} KB")
