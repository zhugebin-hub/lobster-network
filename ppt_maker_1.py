#!/usr/bin/env python3
"""PPT1：科学体育锻炼与身心健康（详细版）"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# 配色
BG_DARK = RGBColor(0x1A, 0x23, 0x7E)
BG_LIGHT = RGBColor(0xF5, 0xF7, 0xFA)
ACCENT = RGBColor(0x00, 0xC8, 0x53)
ACCENT2 = RGBColor(0xFF, 0x6D, 0x00)
ACCENT3 = RGBColor(0x29, 0x62, 0xFF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x21, 0x21, 0x21)
GRAY = RGBColor(0x75, 0x75, 0x75)
LIGHT_GRAY = RGBColor(0xE0, 0xE0, 0xE0)
RED = RGBColor(0xE6, 0x39, 0x46)

def bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color

def shape(slide, l, t, w, h, color):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()
    return s

def rect(slide, l, t, w, h, color):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()
    return s

def tb(slide, l, t, w, h, text, sz=18, color=DARK, bold=False, align=PP_ALIGN.LEFT):
    tx = slide.shapes.add_textbox(l, t, w, h)
    tf = tx.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(sz); p.font.color.rgb = color; p.font.bold = bold
    p.font.name = '微软雅黑'; p.alignment = align
    return tx

def tb_multi(slide, l, t, w, h, lines, sz=16, color=DARK, spacing=6):
    tx = slide.shapes.add_textbox(l, t, w, h)
    tf = tx.text_frame; tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line; p.font.size = Pt(sz); p.font.color.rgb = color
        p.font.name = '微软雅黑'; p.space_after = Pt(spacing)
    return tx

# ====== 第1页：封面 ======
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, BG_DARK)
rect(s, Inches(0), Inches(3.2), Inches(13.333), Inches(0.08), ACCENT)
rect(s, Inches(0), Inches(4.3), Inches(13.333), Inches(0.04), ACCENT2)
tb(s, Inches(1), Inches(1.5), Inches(11.333), Inches(1.5),
   "🏃 科学体育锻炼与身心健康", 44, WHITE, True, PP_ALIGN.CENTER)
tb(s, Inches(2), Inches(3.6), Inches(9.333), Inches(0.6),
   "运动改变生活 · 健康成就未来", 22, ACCENT, False, PP_ALIGN.CENTER)
tb(s, Inches(3), Inches(5.0), Inches(7.333), Inches(0.5),
   "主讲人：______    日期：2026年5月", 16, LIGHT_GRAY, False, PP_ALIGN.CENTER)

# ====== 第2页：目录 ======
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, BG_LIGHT)
tb(s, Inches(0.8), Inches(0.4), Inches(5), Inches(0.8), "📋 目录", 32, BG_DARK, True)
rect(s, Inches(0.8), Inches(1.1), Inches(3), Inches(0.05), ACCENT)
toc = [
    ("01","课程导入与互动调查",ACCENT),("02","体育锻炼对身心健康的影响",ACCENT2),
    ("03","科学锻炼的基本原则",ACCENT3),("04","常见运动方式与适用人群",ACCENT),
    ("05","运动损伤预防与处理",ACCENT2),("06","心理健康与运动处方",ACCENT3),
    ("07","互动环节：问答+讨论+测试",ACCENT),("08","课程总结",ACCENT2),
]
for i,(n,t,c) in enumerate(toc):
    y = Inches(1.5)+Inches(i*0.75)
    shape(s, Inches(1), y, Inches(0.8), Inches(0.55), c)
    tb(s, Inches(1), y+Inches(0.05), Inches(0.8), Inches(0.45), n, 20, WHITE, True, PP_ALIGN.CENTER)
    tb(s, Inches(2.2), y+Inches(0.08), Inches(8), Inches(0.45), t, 18, DARK)

# ====== 第3页：课程导入 ======
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, BG_LIGHT)
tb(s, Inches(0.8), Inches(0.3), Inches(10), Inches(0.8), "01  课程导入", 30, BG_DARK, True)
rect(s, Inches(0.8), Inches(1.0), Inches(4), Inches(0.05), ACCENT)
tb(s, Inches(1.2), Inches(1.4), Inches(11), Inches(0.6),
   "💬 开场思考", 22, ACCENT, True)
tb_multi(s, Inches(1.5), Inches(2.0), Inches(10), Inches(1.5), [
    'WHO（世界卫生组织）指出："运动是人类最好的 medicine"',
    "但你是否知道，错误的运动方式不仅不能带来健康，反而可能造成伤害？",
    "今天我们就来聊聊——如何科学地运动，让身体和心理都受益。",
], 17, DARK, 8)
tb(s, Inches(1.2), Inches(3.8), Inches(11), Inches(0.5),
   "📊 现场互动调查", 22, ACCENT2, True)
tb_multi(s, Inches(1.5), Inches(4.3), Inches(10), Inches(1.5), [
    "1. 你平时喜欢什么运动？坚持了多久？有什么感受？（3位同学分享）",
    "2. 每周运动≥3次的同学请举手？每次≥30分钟的请举手？",
    "3. 你是否有过运动后心情变好的体验？",
], 16, DARK, 8)
tb(s, Inches(1.2), Inches(6.0), Inches(11), Inches(0.5),
   "⏱️ 本环节约5分钟", 14, GRAY, False, PP_ALIGN.CENTER)

# ====== 第4页：身体健康益处 ======
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, BG_LIGHT)
tb(s, Inches(0.8), Inches(0.3), Inches(10), Inches(0.8), "02  体育锻炼对身心健康的影响", 30, BG_DARK, True)
rect(s, Inches(0.8), Inches(1.0), Inches(4), Inches(0.05), ACCENT)
tb(s, Inches(1), Inches(1.4), Inches(5), Inches(0.5), "💪 身体健康益处", 22, ACCENT, True)
body = [
    ("心血管系统","增强心肌力量，降低静息心率","长期有氧运动可使心率下降10-20次/分钟"),
    ("呼吸系统","提高肺活量，增强气体交换效率","运动者肺活量比普通人高20%-30%"),
    ("骨骼肌肉","增加骨密度，预防骨质疏松","负重运动可增加骨密度1%-3%/年"),
    ("代谢系统","改善胰岛素敏感性，控制血糖","运动后24-48小时血糖控制能力提升"),
    ("免疫系统","增强免疫力，减少感冒频率","适度运动者感冒次数减少30%-40%"),
]
for i,(sys,eff,data) in enumerate(body):
    y = Inches(2.0)+Inches(i*0.95)
    shape(s, Inches(1), y, Inches(5), Inches(0.85), WHITE)
    tb(s, Inches(1.2), y+Inches(0.05), Inches(1.5), Inches(0.35), f"🔹 {sys}", 15, BG_DARK, True)
    tb(s, Inches(1.2), y+Inches(0.35), Inches(4.5), Inches(0.25), eff, 12, GRAY)
    tb(s, Inches(3.2), y+Inches(0.35), Inches(2.5), Inches(0.25), data, 11, ACCENT)
# 右侧关键数据
shape(s, Inches(6.8), Inches(1.4), Inches(5.5), Inches(5.5), WHITE)
tb(s, Inches(7.2), Inches(1.6), Inches(4.5), Inches(0.5), "📌 关键数据", 22, ACCENT2, True)
tb_multi(s, Inches(7.2), Inches(2.2), Inches(4.8), Inches(4), [
    "WHO建议：",
    "成年人每周至少150分钟中等强度有氧运动",
    "或75分钟高强度有氧运动",
    "",
    "久坐危害：",
    "每天久坐>8小时，死亡率增加90%",
    "",
    "💡 对比：",
    "• 运动者 vs 久坐者",
    "• 心血管疾病风险降低30%",
    "• 全因死亡率降低20-30%",
    "• 预期寿命延长3-7年",
], 15, DARK, 6)

# ====== 第5页：心理健康益处 ======
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, BG_LIGHT)
tb(s, Inches(0.8), Inches(0.3), Inches(10), Inches(0.8), "02  体育锻炼对身心健康的影响", 30, BG_DARK, True)
rect(s, Inches(0.8), Inches(1.0), Inches(4), Inches(0.05), ACCENT2)
tb(s, Inches(1), Inches(1.4), Inches(5), Inches(0.5), "🧠 心理健康益处", 22, ACCENT2, True)
tb(s, Inches(1), Inches(1.9), Inches(5), Inches(0.4), "神经递质调节", 18, BG_DARK, True)
tb_multi(s, Inches(1.2), Inches(2.3), Inches(4.8), Inches(2), [
    "• 内啡肽释放 → 天然"快乐激素"，产生愉悦感",
    "• 多巴胺分泌 → 提升动力和专注力",
    "• 血清素增加 → 改善情绪，缓解抑郁",
], 14, DARK, 6)
tb(s, Inches(1), Inches(3.8), Inches(5), Inches(0.4), "压力缓解机制", 18, BG_DARK, True)
tb_multi(s, Inches(1.2), Inches(4.2), Inches(4.8), Inches(2), [
    "• 降低皮质醇（压力激素）水平",
    "• 激活副交感神经，促进放松",
    "• 提供"心理断联"时间，暂时脱离压力源",
], 14, DARK, 6)
# 右侧
shape(s, Inches(6.8), Inches(1.4), Inches(5.5), Inches(5.5), WHITE)
tb(s, Inches(7.2), Inches(1.6), Inches(4.5), Inches(0.5), "💬 课堂讨论（5分钟）", 20, ACCENT2, True)
tb_multi(s, Inches(7.2), Inches(2.2), Inches(4.8), Inches(4), [
    "讨论话题：",
    "你是否有过运动后心情变好的体验？",
    "请分享一次具体的经历。",
    "",
    "引导方向：",
    "• 运动前后的情绪变化",
    "• 运动对学习效率的影响",
    "• 运动对睡眠质量的改善",
    "",
    "📝 每组选一名代表分享（2分钟）",
], 15, DARK, 6)

# ====== 第6页：科学锻炼五大原则 ======
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, BG_LIGHT)
tb(s, Inches(0.8), Inches(0.3), Inches(10), Inches(0.8), "03  科学锻炼的基本原则", 30, BG_DARK, True)
rect(s, Inches(0.8), Inches(1.0), Inches(4), Inches(0.05), ACCENT3)
principles = [
    ("🎯","循序渐进","运动量由小到大，强度由低到高","10%原则：每周运动量增加不超过上周的10%","错误：平时不运动，周末跑半马 → 横纹肌溶解"),
    ("⚖️","个性化","根据个人年龄、体质、兴趣选择","考虑时间、场地、健康状况","喜欢才能坚持"),
    ("🔄","全面性","有氧+力量+柔韧+平衡全面发展","FITT原则：频率/强度/时间/类型","每周3-5次，每次30-60分钟"),
    ("⏰","规律性","固定时间、固定项目、固定强度","养成习惯需要21-66天","不规律运动效果大打折扣"),
    ("🛡️","安全性","运动前评估健康状况","充分热身10-15分钟","出现胸痛、头晕立即停止"),
]
for i,(icon,title,desc,tip,warn) in enumerate(principles):
    y = Inches(1.4)+Inches(i*1.1)
    shape(s, Inches(0.8), y, Inches(11.5), Inches(0.95), WHITE)
    tb(s, Inches(1.2), y+Inches(0.05), Inches(1), Inches(0.7), icon, 28, align=PP_ALIGN.CENTER)
    tb(s, Inches(2.5), y+Inches(0.05), Inches(2.5), Inches(0.4), title, 18, BG_DARK, True)
    tb(s, Inches(2.5), y+Inches(0.5), Inches(2.5), Inches(0.4), desc, 13, GRAY)
    tb(s, Inches(5.5), y+Inches(0.1), Inches(3), Inches(0.35), f"💡 {tip}", 13, ACCENT3)
    tb(s, Inches(9), y+Inches(0.1), Inches(3), Inches(0.35), f"⚠️ {warn}", 13, RED)

# ====== 第7页：靶心率与强度评估 ======
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, BG_LIGHT)
tb(s, Inches(0.8), Inches(0.3), Inches(10), Inches(0.8), "03  运动强度的科学评估", 30, BG_DARK, True)
rect(s, Inches(0.8), Inches(1.0), Inches(4), Inches(0.05), ACCENT3)
# 左侧公式
shape(s, Inches(0.8), Inches(1.5), Inches(5.5), Inches(5.2), WHITE)
tb(s, Inches(1.2), Inches(1.7), Inches(4.5), Inches(0.5), "📐 靶心率计算公式", 20, ACCENT3, True)
tb_multi(s, Inches(1.2), Inches(2.3), Inches(4.8), Inches(3), [
    "最大心率 = 220 - 年龄",
    "",
    "靶心率区间 = 最大心率 × (60% ~ 80%)",
    "",
    "举例：20岁学生",
    "最大心率 = 220 - 20 = 200次/分钟",
    "靶心率 = 120 ~ 160次/分钟",
], 16, DARK, 8)
tb(s, Inches(1.2), Inches(5.2), Inches(4.8), Inches(0.5),
   "🔴 现场小测试：自测安静心率（摸桡动脉15秒×4）", 14, RED, True)
# 右侧RPE量表
shape(s, Inches(6.8), Inches(1.5), Inches(5.5), Inches(5.2), WHITE)
tb(s, Inches(7.2), Inches(1.7), Inches(4.5), Inches(0.5), "📊 主观疲劳感觉量表（RPE）", 18, ACCENT3, True)
rpe = [
    ("6-8","非常轻松","低强度"),("9-11","轻松","低-中强度"),
    ("12-14","有些吃力","中等强度 ✅"),("15-17","吃力","高强度"),
    ("18-20","非常吃力","极限强度"),
]
for i,(level,feel,intensity) in enumerate(rpe):
    y = Inches(2.4)+Inches(i*0.85)
    shape(s, Inches(7.2), y, Inches(4.8), Inches(0.7), WHITE if i!=2 else RGBColor(0xE8,0xF5,0xE9))
    tb(s, Inches(7.4), y+Inches(0.05), Inches(1.2), Inches(0.3), level, 16, BG_DARK, True)
    tb(s, Inches(8.6), y+Inches(0.05), Inches(1.5), Inches(0.3), feel, 14, DARK)
    tb(s, Inches(10.2), y+Inches(0.05), Inches(1.5), Inches(0.3), intensity, 13, ACCENT if i==2 else GRAY)

# ====== 第8页：常见运动方式 ======
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, BG_LIGHT)
tb(s, Inches(0.8), Inches(0.3), Inches(10), Inches(0.8), "04  常见运动方式与适用人群", 30, BG_DARK, True)
rect(s, Inches(0.8), Inches(1.0), Inches(4), Inches(0.05), ACCENT)
exercises = [
    ("🏃","有氧运动","跑步、游泳、骑行、快走","增强心肺、减脂","大多数人群","膝关节不好选游泳"),
    ("💪","力量训练","举重、器械、自重训练","增肌、提高代谢","青少年、中老年","动作规范第一"),
    ("🧘","柔韧训练","瑜伽、拉伸、太极","提高柔韧性、放松","久坐人群、压力大","不要弹震式拉伸"),
    ("⚽","球类运动","篮球、足球、羽毛球","综合体能、社交","喜欢团队活动者","注意扭伤风险"),
    ("🔥","HIIT","间歇训练、Tabata","高效燃脂、节省时间","有一定基础者","不适合初学者"),
]
for i,(icon,name,items,eff,who,note) in enumerate(exercises):
    y = Inches(1.4)+Inches(i*1.1)
    shape(s, Inches(0.8), y, Inches(11.5), Inches(0.95), WHITE)
    tb(s, Inches(1.2), y+Inches(0.05), Inches(1.5), Inches(0.35), f"{icon} {name}", 17, BG_DARK, True)
    tb(s, Inches(3), y+Inches(0.05), Inches(3), Inches(0.35), items, 13, GRAY)
    tb(s, Inches(6.5), y+Inches(0.05), Inches(2.5), Inches(0.35), f"功效：{eff}", 13, ACCENT)
    tb(s, Inches(9.5), y+Inches(0.05), Inches(2.5), Inches(0.35), f"⚠️ {note}", 12, RED)

# 右侧人群推荐
shape(s, Inches(0.8), Inches(5.0), Inches(11.5), Inches(2.2), WHITE)
tb(s, Inches(1.2), Inches(5.1), Inches(10), Inches(0.4), "👥 不同人群的推荐运动方案", 18, BG_DARK, True)
tb_multi(s, Inches(1.2), Inches(5.5), Inches(10.5), Inches(1.5), [
    "🎓 大学生（18-25岁）：每周4-5次，每次40-60分钟，有氧+力量结合 → 跑步、打球、健身",
    "💼 久坐上班族：每周3-4次，每次30-45分钟，有氧+拉伸为主 → 快走、游泳、瑜伽",
    "👴 中老年人（50岁以上）：每周3-5次，每次30-40分钟，低冲击有氧+平衡 → 散步、太极、水中运动",
], 14, DARK, 5)

# ====== 第9页：运动损伤预防 ======
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, BG_LIGHT)
tb(s, Inches(0.8), Inches(0.3), Inches(10), Inches(0.8), "05  运动损伤预防与处理", 30, BG_DARK, True)
rect(s, Inches(0.8), Inches(1.0), Inches(4), Inches(0.05), ACCENT2)
# 左侧常见损伤
shape(s, Inches(0.8), Inches(1.5), Inches(5.8), Inches(5.3), WHITE)
tb(s, Inches(1.2), Inches(1.7), Inches(4.5), Inches(0.5), "🛡️ 常见运动损伤", 20, ACCENT2, True)
injuries = [
    ("肌肉拉伤","大腿后侧、小腿","热身不足、过度拉伸","充分热身、循序渐进"),
    ("韧带扭伤","踝关节、膝关节","落地不稳、急转急停","穿合适鞋子、加强肌肉力量"),
    ("应力性骨折","胫骨、足部","过度使用、突然增加运动量","遵循10%原则"),
    ("肩袖损伤","肩部","重复性overhead动作","加强肩袖肌群训练"),
]
for i,(name,part,cause,prevent) in enumerate(injuries):
    y = Inches(2.3)+Inches(i*1.1)
    tb(s, Inches(1.2), y+Inches(0.05), Inches(2), Inches(0.3), f"🔸 {name}", 15, BG_DARK, True)
    tb(s, Inches(1.2), y+Inches(0.35), Inches(1.5), Inches(0.25), f"部位：{part}", 12, GRAY)
    tb(s, Inches(3), y+Inches(0.35), Inches(3.2), Inches(0.25), f"原因：{cause}", 12, RED)
    tb(s, Inches(1.2), y+Inches(0.65), Inches(4.8), Inches(0.25), f"✅ 预防：{prevent}", 12, ACCENT)
# 右侧RICE
shape(s, Inches(7), Inches(1.5), Inches(5.5), Inches(5.3), WHITE)
tb(s, Inches(7.4), Inches(1.7), Inches(4.5), Inches(0.5), "🏥 RICE处理原则", 20, ACCENT2, True)
tb(s, Inches(7.4), Inches(2.2), Inches(4.5), Inches(0.3), "急性损伤发生后48小时内：", 14, GRAY)
rice = [("R","Rest","休息","立即停止运动，避免患肢承重"),
        ("I","Ice","冰敷","每次15-20分钟，每2-3小时一次"),
        ("C","Compression","加压","弹性绷带包扎，从远端向近端"),
        ("E","Elevation","抬高","患肢抬高至心脏水平以上")]
for i,(letter,en,cn,desc) in enumerate(rice):
    y = Inches(2.6)+Inches(i*1.0)
    shape(s, Inches(7.4), y, Inches(0.6), Inches(0.8), ACCENT2)
    tb(s, Inches(7.4), y+Inches(0.1), Inches(0.6), Inches(0.5), letter, 24, WHITE, True, PP_ALIGN.CENTER)
    tb(s, Inches(8.2), y+Inches(0.05), Inches(2), Inches(0.3), f"{en}（{cn}）", 15, BG_DARK, True)
    tb(s, Inches(8.2), y+Inches(0.4), Inches(3.8), Inches(0.3), desc, 12, GRAY)
tb(s, Inches(7.4), Inches(6.2), Inches(4.5), Inches(0.4),
   "⚠️ 严重情况及时就医！", 14, RED, True)

# ====== 第10页：心理健康与运动处方 ======
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, BG_LIGHT)
tb(s, Inches(0.8), Inches(0.3), Inches(10), Inches(0.8), "06  心理健康与运动处方", 30, BG_DARK, True)
rect(s, Inches(0.8), Inches(1.0), Inches(4), Inches(0.05), ACCENT2)
prescriptions = [
    ("😰","缓解焦虑","慢跑+瑜伽","每周3次，每次30-40分钟","降低皮质醇，增加内啡肽"),
    ("😔","改善抑郁","团体运动+户外","每周2-3次，每次40分钟","社交互动+阳光+运动三重作用"),
    ("😴","改善睡眠","晚饭后散步+拉伸","每天傍晚，30分钟","调节生物钟，促进褪黑素分泌"),
    ("😤","释放压力","HIIT/拳击/跑步","每周2次，每次20-30分钟","高强度释放压力，产生成就感"),
    ("🤯","提升专注","游泳/跳绳/舞蹈","每周3次，每次30分钟","提高脑部血流量，促进BDNF分泌"),
]
for i,(emoji,issue,solution,freq,mechanism) in enumerate(prescriptions):
    y = Inches(1.4)+Inches(i*1.1)
    shape(s, Inches(0.8), y, Inches(11.5), Inches(0.95), WHITE)
    tb(s, Inches(1.2), y+Inches(0.05), Inches(2), Inches(0.35), f"{emoji} {issue}", 17, BG_DARK, True)
    tb(s, Inches(3.5), y+Inches(0.05), Inches(2.5), Inches(0.35), f"推荐：{solution}", 14, ACCENT2)
    tb(s, Inches(6.5), y+Inches(0.05), Inches(2.5), Inches(0.35), f"频率：{freq}", 13, GRAY)
    tb(s, Inches(9.5), y+Inches(0.05), Inches(2.5), Inches(0.35), f"机制：{mechanism}", 12, ACCENT3)

# 右侧补充
shape(s, Inches(0.8), Inches(5.0), Inches(11.5), Inches(2.2), WHITE)
tb(s, Inches(1.2), Inches(5.1), Inches(10), Inches(0.4), "🌙 运动与睡眠质量", 18, BG_DARK, True)
tb_multi(s, Inches(1.2), Inches(5.5), Inches(10.5), Inches(1.5), [
    "运动改善睡眠的机制：升高体温后下降→触发睡眠信号；消耗能量→增加睡眠驱动力；缓解焦虑→减少入睡困难",
    "✅ 最佳运动时间：下午4-6点（体温最高）/ 早晨（调节生物钟）",
    "⚠️ 睡前2小时内避免高强度运动",
    "",
    "🚨 运动成瘾的信号：即使受伤也坚持运动 / 不运动就焦虑内疚 / 运动成为唯一减压方式",
], 14, DARK, 5)

# ====== 第11页：互动环节-知识问答 ======
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, BG_DARK)
rect(s, Inches(0), Inches(3.5), Inches(13.333), Inches(0.06), ACCENT)
tb(s, Inches(1), Inches(0.5), Inches(11.333), Inches(1),
   "🎯 互动环节：知识问答", 36, WHITE, True, PP_ALIGN.CENTER)
tb(s, Inches(2), Inches(1.5), Inches(9.333), Inches(0.5),
   "答对有小惊喜哦！🎁", 20, ACCENT, False, PP_ALIGN.CENTER)
questions = [
    "Q1：WHO建议成年人每周至少进行多少分钟的中等强度有氧运动？",
    "    A. 60分钟   B. 150分钟   C. 300分钟   D. 90分钟",
    "    ✅ 答案：B（150分钟）",
    "",
    "Q2：运动后肌肉酸痛（DOMS）一般持续多久？",
    "    A. 1-2天    B. 3-5天    C. 一周      D. 两周以上",
    "    ✅ 答案：B（3-5天，通常在24-72小时达到高峰）",
    "",
    "Q3：以下哪种运动最适合缓解焦虑？",
    "    A. 短跑冲刺  B. 瑜伽冥想  C. 举重训练  D. 跳绳",
    "    ✅ 答案：B（瑜伽冥想，有节奏的深呼吸配合柔和运动最能降低焦虑）",
    "",
    "Q4：20岁学生的靶心率区间大约是多少？",
    "    ✅ 答案：120-160次/分钟（最大心率200×60%~80%）",
]
tb(s, Inches(1.5), Inches(2.2), Inches(10.333), Inches(4.5),
   "\n".join(questions), 17, LIGHT_GRAY, False, PP_ALIGN.LEFT)

# ====== 第12页：互动环节-小组讨论 ======
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, BG_LIGHT)
tb(s, Inches(0.8), Inches(0.3), Inches(10), Inches(0.8), "💬 小组讨论环节", 30, BG_DARK, True)
rect(s, Inches(0.8), Inches(1.0), Inches(4), Inches(0.05), ACCENT3)
topics = [
    ("讨论主题 1","你最喜欢的运动是什么？它给你的生活带来了哪些改变？","每组分享2分钟",ACCENT),
    ("讨论主题 2","大学生应该如何平衡学业和体育锻炼？","引导：碎片化运动、课间拉伸、走路去上课",ACCENT2),
    ("讨论主题 3","如何克服"不想运动"的心理障碍？","引导：找运动伙伴、设定小目标、选择喜欢的运动",ACCENT3),
]
for i,(label,topic,guide,color) in enumerate(topics):
    y = Inches(1.5)+Inches(i*1.8)
    shape(s, Inches(0.8), y, Inches(11.5), Inches(1.5), WHITE)
    rect(s, Inches(0.8), y, Inches(0.12), Inches(1.5), color)
    tb(s, Inches(1.3), y+Inches(0.15), Inches(3), Inches(0.4), label, 20, color, True)
    tb(s, Inches(1.3), y+Inches(0.6), Inches(10), Inches(0.5), topic, 18, DARK)
    tb(s, Inches(1.3), y+Inches(1.1), Inches(10), Inches(0.3), f"📝 {guide}", 14, GRAY)
tb(s, Inches(1), Inches(6.5), Inches(11), Inches(0.5),
   "📝 每组选一名代表分享讨论结果（每组3分钟）", 16, GRAY, False, PP_ALIGN.CENTER)

# ====== 第13页：互动环节-现场测试 ======
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, BG_LIGHT)
tb(s, Inches(0.8), Inches(0.3), Inches(10), Inches(0.8), "🏋️ 现场体能小测试", 30, BG_DARK, True)
rect(s, Inches(0.8), Inches(1.0), Inches(4), Inches(0.05), ACCENT2)
tb(s, Inches(1.2), Inches(1.4), Inches(10), Inches(0.4),
   "分组进行，互相监督记录", 18, GRAY, False, PP_ALIGN.CENTER)
tests = [
    ("平板支撑","核心力量","男生≥45秒 / 女生≥30秒","优秀：≥90秒","💪"),
    ("坐位体前屈","柔韧性","男生≥8cm / 女生≥12cm","优秀：≥18cm","🧘"),
    ("1分钟跳绳","协调性","≥120次","优秀：≥160次","🤸"),
    ("闭眼单脚站立","平衡能力","男生≥25秒 / 女生≥30秒","优秀：≥45秒","🦩"),
]
for i,(name,desc,standard,excellent,icon) in enumerate(tests):
    col = i%2; row = i//2
    x = Inches(0.8)+Inches(col*6); y = Inches(2.0)+Inches(row*1.8)
    shape(s, x, y, Inches(5.5), Inches(1.5), WHITE)
    tb(s, x+Inches(0.3), y+Inches(0.1), Inches(2), Inches(0.4), f"{icon} {name}", 20, BG_DARK, True)
    tb(s, x+Inches(0.3), y+Inches(0.5), Inches(4.5), Inches(0.3), desc, 14, GRAY)
    tb(s, x+Inches(0.3), y+Inches(0.8), Inches(4.5), Inches(0.3), f"📏 达标：{standard}", 14, ACCENT2)
    tb(s, x+Inches(0.3), y+Inches(1.1), Inches(4.5), Inches(0.3), f"⭐ 优秀：{excellent}", 13, ACCENT3)
tb(s, Inches(1), Inches(6.2), Inches(11), Inches(0.8),
   "💬 测试后讨论：你的优势项目是什么？薄弱环节是什么？根据测试结果制定4周改善计划",
   15, GRAY, False, PP_ALIGN.CENTER)

# ====== 第14页：总结 ======
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, BG_DARK)
rect(s, Inches(0), Inches(3.5), Inches(13.333), Inches(0.06), ACCENT)
tb(s, Inches(1), Inches(0.5), Inches(11.333), Inches(1),
   "📌 课程总结", 36, WHITE, True, PP_ALIGN.CENTER)
summary = [
    "✅ 体育锻炼是身心健康的基石",
    "✅ 科学锻炼要遵循五大原则",
    "✅ 选择合适的运动方式很重要",
    "✅ 预防运动损伤不可忽视",
    "✅ 运动是最好的"心理处方"",
]
tb(s, Inches(2), Inches(2), Inches(9.333), Inches(3.5),
   "\n\n".join(summary), 22, LIGHT_GRAY, False, PP_ALIGN.CENTER)
tb(s, Inches(2), Inches(5.5), Inches(9.333), Inches(0.8),
   "🏃‍♂️ 每天锻炼一小时，健康工作五十年，幸福生活一辈子！",
   20, ACCENT, True, PP_ALIGN.CENTER)

# ====== 第15页：结束页 ======
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, BG_DARK)
rect(s, Inches(0), Inches(3.2), Inches(13.333), Inches(0.08), ACCENT)
rect(s, Inches(0), Inches(4.3), Inches(13.333), Inches(0.04), ACCENT2)
tb(s, Inches(1), Inches(1.5), Inches(11.333), Inches(1.5),
   "🙏 感谢聆听！", 44, WHITE, True, PP_ALIGN.CENTER)
tb(s, Inches(2), Inches(3.6), Inches(9.333), Inches(0.6),
   "让我们一起动起来，拥抱健康生活！", 22, ACCENT, False, PP_ALIGN.CENTER)
tb(s, Inches(3), Inches(5.0), Inches(7.333), Inches(0.5),
   "Q & A    欢迎提问与交流", 18, LIGHT_GRAY, False, PP_ALIGN.CENTER)

output = "/home/admin/.openclaw/workspace/科学体育锻炼与身心健康-详细版.pptx"
prs.save(output)
print(f"✅ PPT1 已保存: {output}")
