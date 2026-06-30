#!/usr/bin/env python3
# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# 创建演示文稿 - 宽屏 16:9
prs = Presentation()
prs.slide_width = Cm(33.87)  # 13.33 英寸
prs.slide_height = Cm(19.05)  # 7.5 英寸

# 定义配色方案
colors = {
    'primary': RGBColor(102, 187, 106),      # 清新绿
    'primary_dark': RGBColor(76, 161, 80),   # 深绿
    'secondary': RGBColor(66, 165, 245),     # 天空蓝
    'accent': RGBColor(255, 167, 38),        # 活力橙
    'warm': RGBColor(255, 204, 128),         # 暖黄
    'text_dark': RGBColor(51, 51, 51),       # 深灰文字
    'text_light': RGBColor(102, 102, 102),   # 浅灰文字
    'white': RGBColor(255, 255, 255),        # 白色
    'bg_light': RGBColor(245, 250, 245),     # 浅绿背景
}

def add_gradient_bg(slide, color1, color2):
    """添加渐变背景效果（用矩形模拟）"""
    # 上半部分
    shape1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Cm(9.5))
    shape1.fill.solid()
    shape1.fill.fore_color.rgb = color1
    shape1.line.fill.background()
    # 下半部分
    shape2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Cm(9.5), prs.slide_width, Cm(9.55))
    shape2.fill.solid()
    shape2.fill.fore_color.rgb = color2
    shape2.line.fill.background()

def add_decorative_circle(slide, x, y, size, color, transparency=0.3):
    """添加装饰圆形"""
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, size, size)
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.fill.transparency = transparency
    circle.line.fill.background()

def add_title_slide(prs, title, subtitle):
    """精美的标题页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 渐变背景
    add_gradient_bg(slide, colors['bg_light'], RGBColor(255, 255, 255))
    
    # 装饰元素 - 左上角大圆
    add_decorative_circle(slide, Cm(-3), Cm(-3), Cm(8), colors['primary'], 0.1)
    add_decorative_circle(slide, Cm(-1), Cm(-1), Cm(5), colors['secondary'], 0.15)
    
    # 装饰元素 - 右下角
    add_decorative_circle(slide, Cm(28), Cm(14), Cm(6), colors['accent'], 0.1)
    add_decorative_circle(slide, Cm(30), Cm(16), Cm(4), colors['warm'], 0.15)
    
    # 主标题框
    title_box = slide.shapes.add_textbox(Cm(2), Cm(6), Cm(29.87), Cm(5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = colors['primary_dark']
    p.alignment = PP_ALIGN.CENTER
    p.font.name = 'Microsoft YaHei'
    
    # 副标题
    sub_box = slide.shapes.add_textbox(Cm(2), Cm(10.5), Cm(29.87), Cm(3))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(28)
    p.font.color.rgb = colors['text_light']
    p.alignment = PP_ALIGN.CENTER
    p.font.name = 'Microsoft YaHei'
    
    # 底部装饰线
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(5), Cm(17.5), Cm(23.87), Cm(0.15))
    line.fill.solid()
    line.fill.fore_color.rgb = colors['primary']
    line.line.fill.background()
    
    return slide

def add_content_slide(prs, title, items, icon_colors=None):
    """精美的内容页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 顶部色块
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Cm(2.5))
    header.fill.solid()
    header.fill.fore_color.rgb = colors['primary']
    header.line.fill.background()
    
    # 标题
    title_box = slide.shapes.add_textbox(Cm(1), Cm(0.5), Cm(31.87), Cm(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = colors['white']
    p.alignment = PP_ALIGN.CENTER
    p.font.name = 'Microsoft YaHei'
    
    # 内容区域背景
    bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(1), Cm(3), Cm(31.87), Cm(14))
    bg.fill.solid()
    bg.fill.fore_color.rgb = colors['bg_light']
    bg.line.color.rgb = colors['primary']
    bg.line.width = Pt(2)
    
    # 内容项
    if icon_colors is None:
        icon_colors = [colors['primary'], colors['secondary'], colors['accent'], colors['warm']]
    
    for i, item in enumerate(items):
        y_pos = Cm(4 + i * 2.5)
        
        # 图标背景圆
        if isinstance(item, tuple) and len(item) > 1:
            icon = item[0]
            text = item[1]
            circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Cm(2), y_pos, Cm(1.2), Cm(1.2))
            circle.fill.solid()
            circle.fill.fore_color.rgb = icon_colors[i % len(icon_colors)]
            circle.line.fill.background()
            
            # 图标文字（简化为色块）
            icon_box = slide.shapes.add_textbox(Cm(2.3), y_pos + Cm(0.3), Cm(0.6), Cm(0.6))
            tf = icon_box.text_frame
            p = tf.paragraphs[0]
            p.text = icon if icon else '●'
            p.font.size = Pt(20)
            p.font.color.rgb = colors['white']
            p.alignment = PP_ALIGN.CENTER
            
            # 内容文字
            text_box = slide.shapes.add_textbox(Cm(4), y_pos, Cm(28), Cm(1.5))
            tf = text_box.text_frame
            p = tf.paragraphs[0]
            p.text = text
            p.font.size = Pt(24)
            p.font.color.rgb = colors['text_dark']
            p.font.name = 'Microsoft YaHei'
        else:
            text_box = slide.shapes.add_textbox(Cm(2), y_pos, Cm(30), Cm(1.5))
            tf = text_box.text_frame
            p = tf.paragraphs[0]
            p.text = item
            p.font.size = Pt(24)
            p.font.color.rgb = colors['text_dark']
            p.font.name = 'Microsoft YaHei'
    
    return slide

def add_feature_slide(prs, title, features):
    """三大功能页 - 卡片式布局"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 顶部色块
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Cm(2.5))
    header.fill.solid()
    header.fill.fore_color.rgb = colors['primary']
    header.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Cm(1), Cm(0.5), Cm(31.87), Cm(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = colors['white']
    p.alignment = PP_ALIGN.CENTER
    p.font.name = 'Microsoft YaHei'
    
    # 三个功能卡片
    card_width = Cm(10)
    card_height = Cm(11)
    card_y = Cm(4)
    spacing = Cm(0.5)
    start_x = Cm(2)
    
    card_colors = [colors['primary'], colors['secondary'], colors['accent']]
    
    for i, feature in enumerate(features):
        x_pos = start_x + i * (card_width + spacing)
        
        # 卡片背景
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x_pos, card_y, card_width, card_height)
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(255, 255, 255)
        card.line.color.rgb = card_colors[i]
        card.line.width = Pt(3)
        
        # 卡片顶部色条
        top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x_pos, card_y, card_width, Cm(1.5))
        top_bar.fill.solid()
        top_bar.fill.fore_color.rgb = card_colors[i]
        top_bar.line.fill.background()
        
        # 功能标题
        title_box = slide.shapes.add_textbox(x_pos + Cm(0.5), card_y + Cm(0.3), card_width - Cm(1), Cm(0.8))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = feature['title']
        p.font.size = Pt(26)
        p.font.bold = True
        p.font.color.rgb = colors['white']
        p.alignment = PP_ALIGN.CENTER
        p.font.name = 'Microsoft YaHei'
        
        # 功能描述
        desc_box = slide.shapes.add_textbox(x_pos + Cm(0.5), card_y + Cm(2), card_width - Cm(1), card_height - Cm(2.5))
        tf = desc_box.text_frame
        tf.word_wrap = True
        for j, line in enumerate(feature['desc']):
            if j == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = line
            p.font.size = Pt(18)
            p.font.color.rgb = colors['text_dark']
            p.font.name = 'Microsoft YaHei'
            p.space_after = Pt(8)
    
    return slide

def add_example_slide(prs, title, examples):
    """互动示例页 - 对话气泡风格"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 顶部色块
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Cm(2.5))
    header.fill.solid()
    header.fill.fore_color.rgb = colors['primary']
    header.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Cm(1), Cm(0.5), Cm(31.87), Cm(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = colors['white']
    p.alignment = PP_ALIGN.CENTER
    p.font.name = 'Microsoft YaHei'
    
    # 两个示例左右布局
    for i, example in enumerate(examples):
        x_offset = i * Cm(17)
        
        # 场景标题
        scene_box = slide.shapes.add_textbox(Cm(2) + x_offset, Cm(3), Cm(15), Cm(1))
        tf = scene_box.text_frame
        p = tf.paragraphs[0]
        p.text = example['scene']
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = colors['primary_dark']
        p.font.name = 'Microsoft YaHei'
        
        # 对话背景
        dialog_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(2) + x_offset, Cm(4.2), Cm(15), Cm(11))
        dialog_bg.fill.solid()
        dialog_bg.fill.fore_color.rgb = colors['bg_light']
        dialog_bg.line.color.rgb = colors['primary']
        dialog_bg.line.width = Pt(2)
        
        # 对话内容
        dialog_box = slide.shapes.add_textbox(Cm(2.5) + x_offset, Cm(4.5), Cm(14), Cm(10))
        tf = dialog_box.text_frame
        tf.word_wrap = True
        for j, line in enumerate(example['dialog']):
            if j == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = line
            p.font.size = Pt(16)
            p.font.color.rgb = colors['text_dark']
            p.font.name = 'Microsoft YaHei'
            p.space_after = Pt(10)
    
    return slide

# ============ 创建幻灯片 ============

# 1. 封面页
add_title_slide(prs, "潮潮姐智能体", "海宁市王国维小学教育集团 · AI 心理伙伴")

# 2. 潮潮姐是谁
add_content_slide(prs, "潮潮姐是谁？", [
    ('🤖', "学校专属的 AI 心理健康伙伴"),
    ('🎨', "基于学校吉祥物'潮娃'设计"),
    ('☀️', "形象亲切阳光，像贴心大朋友"),
    ('🛡️', "使命：陪伴孩子，守护心理健康"),
    ('👦', "做校园里随时在线的心理小卫士"),
])

# 3. 使用场景
add_content_slide(prs, "使用场景", [
    ('📍', "部署位置：心理健康教室专用平板"),
    ('⏰', "使用时间：课间、午休、课后随时可用"),
    ('🔒', "隐私保护：一对一平板交互"),
    ('💡', "适合所有孩子，包括内向害羞的孩子"),
    ('👩‍🏫', "定位：心理老师的辅助工具"),
], [colors['secondary'], colors['accent'], colors['warm'], colors['primary'], colors['secondary']])

# 4. 三大核心功能 - 卡片式
add_feature_slide(prs, "三大核心功能", [
    {
        'title': '🌳 情绪树洞',
        'desc': [
            '• 倾诉烦恼、不开心',
            '• 共情安慰',
            '• 情绪支持'
        ]
    },
    {
        'title': '📚 心理小课堂',
        'desc': [
            '• 趣味科普心理知识',
            '• 教情绪调节小技巧',
            '• 寓教于乐'
        ]
    },
    {
        'title': '🌱 成长陪伴',
        'desc': [
            '• 记录成长',
            '• 正向激励',
            '• 提供亲子沟通建议',
            '• 家校协同关注心理健康'
        ]
    },
])

# 5. 互动示例
add_example_slide(prs, "互动示例", [
    {
        'scene': '场景一：考试没考好',
        'dialog': [
            '孩子："潮潮姐姐，我今天考试没考好，很难过"',
            '',
            '潮潮姐："抱抱你，没考好心里肯定不好受吧，我能理解这种感觉"',
            '',
            '潮潮姐："你已经很努力了，一次的成绩不代表全部哦，我们一起来看看下次怎么可以做得更好一点"'
        ]
    },
    {
        'scene': '场景二：上课走神',
        'dialog': [
            '孩子："潮潮姐姐，我上课总是容易走神，注意力不集中，该怎么办呢？"',
            '',
            '潮潮姐："这是很普遍的现象，别担心"',
            '',
            '潮潮姐："推荐番茄学习法——专注听课 20 分钟，休息 5 分钟，像小番茄一样一个一个完成任务，慢慢养成好习惯"'
        ]
    },
])

# 6. 价值与未来
add_content_slide(prs, "价值与未来规划", [
    ('💪', "当前价值"),
    ('', "  • 随时在线的心理支持，不用等老师有空"),
    ('', "  • 减轻心理老师工作压力，提升教育效率"),
    ('', "  • 构建家校协同的心理健康教育新模式"),
    ('', ""),
    ('🚀', "未来规划"),
    ('', "  • 优化对话能力，提供更精准的个性化辅导"),
    ('', "  • 拓展更多互动形式：心理小游戏、小故事"),
    ('', "  • 让孩子在玩中学到心理知识"),
], [colors['primary'], colors['secondary'], colors['secondary'], colors['secondary'], colors['accent'], colors['accent'], colors['accent'], colors['accent']])

# 7. 结束页
add_title_slide(prs, "谢谢大家！", "潮潮姐——孩子身边随时能找到的心理小卫士")

# 保存
prs.save('/home/admin/.openclaw/workspace/潮潮姐智能体演示_精美版.pptx')
print("精美 PPT 生成成功！")
