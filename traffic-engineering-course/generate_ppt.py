#!/usr/bin/env python3
"""
Generate comprehensive Traffic Engineering course PPT
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Create presentation with 16:9 aspect ratio
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color scheme
DARK_BG = RGBColor(10, 22, 40)
LIGHT_BG = RGBColor(240, 244, 248)
ACCENT = RGBColor(79, 195, 247)
ACCENT_SOFT = RGBColor(129, 212, 250)
GREEN = RGBColor(102, 187, 106)
AMBER = RGBColor(255, 183, 77)
RED = RGBColor(239, 83, 80)
PURPLE = RGBColor(171, 71, 188)
TEXT_SOFT = RGBColor(196, 212, 224)
MUTED = RGBColor(90, 122, 148)
WHITE = RGBColor(255, 255, 255)
DARK_SURFACE = RGBColor(15, 32, 53)

def add_shape(slide, shape_type, left, top, width, height, fill_color=None, line_color=None):
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape

def add_text_box(slide, text, left, top, width, height, font_size=18, font_color=WHITE, 
                 bold=False, font_name='Microsoft YaHei', alignment=PP_ALIGN.LEFT, 
                 line_spacing=None, is_italic=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = alignment
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.name = font_name
    p.font.italic = is_italic
    if line_spacing:
        p.font.line_spacing = line_spacing
    return txBox

def add_bullets(slide, items, left, top, width, height, font_size=16, font_color=TEXT_SOFT,
                font_name='Microsoft YaHei', bullet_color=ACCENT, spacing=0):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.name = font_name
        p.level = 0
        if spacing > 0:
            p.space_after = Pt(spacing)
    
    return txBox

def add_section_divider(slide, chapter_num, chapter_title, subtitle, accent_color=ACCENT):
    # Background
    bg = add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5), 
                   fill_color=DARK_BG)
    
    # Grid pattern lines
    for i in range(0, 14, 2):
        add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(i), Inches(0), Inches(0.02), Inches(7.5), 
                  fill_color=RGBColor(79, 195, 247))
    
    # Accent line
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(2.5), Inches(0.1), Inches(2.5),
              fill_color=accent_color)
    
    # Chapter number
    add_text_box(slide, chapter_num, Inches(1.8), Inches(2.5), Inches(11), Inches(0.6), 
                 font_size=20, font_color=MUTED, font_name='Consolas', bold=False)
    
    # Chapter title
    add_text_box(slide, chapter_title, Inches(1.8), Inches(3.2), Inches(11), Inches(1.2), 
                 font_size=56, font_color=WHITE, bold=True)
    
    # Subtitle
    add_text_box(slide, subtitle, Inches(1.8), Inches(4.6), Inches(10), Inches(0.8), 
                 font_size=24, font_color=TEXT_SOFT)

# ============== SLIDE 1: COVER ==============
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

# Background
bg = add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5), 
               fill_color=DARK_BG)

# Decorative road lines
for i in range(8):
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(11 + i * 0.15), Inches(0), Inches(0.04), Inches(7.5),
              fill_color=RGBColor(79, 195, 247))

# Subtitle top
add_text_box(slide, '高等学校 · 交通运输类专业核心课程', Inches(1), Inches(0.8), Inches(11), Inches(0.5),
             font_size=16, font_color=MUTED, font_name='Consolas')

# Main title
add_text_box(slide, '交通工程学', Inches(1), Inches(2.2), Inches(10), Inches(1.5),
             font_size=80, font_color=WHITE, bold=True)

# English subtitle
add_text_box(slide, 'Traffic Engineering', Inches(1), Inches(3.6), Inches(10), Inches(0.8),
             font_size=36, font_color=ACCENT_SOFT)

# Description
add_text_box(slide, '系统学习交通流理论、交通规划、交通管理与控制的核心理论与实践方法', 
             Inches(1), Inches(4.6), Inches(10), Inches(0.8),
             font_size=22, font_color=TEXT_SOFT)

# Bottom keywords
keywords = ['📚 理论体系', ' 分析方法', '🛣️ 工程实践', '🚦 智能交通']
for i, kw in enumerate(keywords):
    add_text_box(slide, kw, Inches(1 + i * 2.8), Inches(6.2), Inches(2.5), Inches(0.5),
                 font_size=16, font_color=MUTED, font_name='Consolas')

# ============== SLIDE 2: COURSE OVERVIEW ==============
slide = prs.slides.add_slide(prs.slide_layouts[6])

bg = add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5),
               fill_color=DARK_BG)

add_text_box(slide, '课程概述', Inches(1), Inches(0.6), Inches(11), Inches(0.5),
             font_size=16, font_color=MUTED, font_name='Consolas')

add_text_box(slide, '为什么要学习交通工程学？', Inches(1), Inches(1.1), Inches(11), Inches(0.8),
             font_size=40, font_color=WHITE, bold=True)

# Four cards
card_data = [
    ('🌆', '城市化挑战', '中国城镇化率已超过 65%，城市交通拥堵、\n停车难、出行效率低下成为制约城市\n发展的瓶颈问题'),
    ('📊', '数据驱动决策', '运用交通调查、数据分析和模型预测，\n为交通规划、设计和管理提供科学依据'),
    ('🤖', '智慧交通发展', '车路协同、自动驾驶、MaaS 等新技术\n正在重塑交通系统，需要新型工程人才'),
    ('🌱', '可持续发展', '双碳目标下，构建绿色、低碳、高效的\n综合交通运输体系是国家战略需求'),
]

for i, (emoji, title, desc) in enumerate(card_data):
    left = 1 + (i % 2) * 5.8
    top = 2.3 + (i // 2) * 2.3
    
    # Card background
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(5.5), Inches(2),
              fill_color=DARK_SURFACE, line_color=RGBColor(79, 195, 247))
    
    add_text_box(slide, emoji, Inches(left + 0.3), Inches(top + 0.15), Inches(1), Inches(0.6),
                 font_size=28)
    add_text_box(slide, title, Inches(left + 1.2), Inches(top + 0.15), Inches(4), Inches(0.6),
                 font_size=20, font_color=WHITE, bold=True)
    add_text_box(slide, desc, Inches(left + 0.3), Inches(top + 0.8), Inches(5), Inches(1.1),
                 font_size=14, font_color=TEXT_SOFT)

# ============== SLIDE 3: TABLE OF CONTENTS ==============
slide = prs.slides.add_slide(prs.slide_layouts[6])

bg = add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5),
               fill_color=DARK_BG)

add_text_box(slide, '课程目录', Inches(1), Inches(0.6), Inches(11), Inches(0.5),
             font_size=16, font_color=MUTED, font_name='Consolas')

add_text_box(slide, '六大核心模块', Inches(1), Inches(1.1), Inches(11), Inches(0.8),
             font_size=40, font_color=WHITE, bold=True)

chapters = [
    ('01', '交通工程学导论', '学科定义、发展历程、研究对象', ACCENT),
    ('02', '交通流理论', '交通流特性、宏观微观模型、通行能力', GREEN),
    ('03', '交通调查与数据分析', '流量速度密度调查、OD 调查、数据方法', AMBER),
    ('04', '道路交通规划', '四阶段法、交通生成与分布、方式划分', PURPLE),
    ('05', '交叉口设计与信号控制', '渠化设计、信号配时、延误分析', RED),
    ('06', '交通管理与智能交通系统', '交通组织、ITS 架构、车路协同', ACCENT_SOFT),
]

for i, (num, title, desc, color) in enumerate(chapters):
    top = 2.3 + i * 0.75
    
    # Left accent bar
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(1), Inches(top), Inches(0.05), Inches(0.65),
              fill_color=color)
    
    # Background
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(1), Inches(top), Inches(11.333), Inches(0.65),
              fill_color=DARK_SURFACE, line_color=RGBColor(79, 195, 247))
    
    add_text_box(slide, num, Inches(1.3), Inches(top + 0.1), Inches(0.8), Inches(0.5),
                 font_size=18, font_color=color, font_name='Consolas', bold=True)
    add_text_box(slide, title, Inches(2.1), Inches(top + 0.08), Inches(5), Inches(0.5),
                 font_size=20, font_color=WHITE, bold=True)
    add_text_box(slide, desc, Inches(7.2), Inches(top + 0.1), Inches(5), Inches(0.5),
                 font_size=14, font_color=MUTED)

# ============== SLIDE 4: CHAPTER 1 - INTRODUCTION ==============
slide = prs.slides.add_slide(prs.slide_layouts[6])

bg = add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5),
               fill_color=DARK_BG)

add_text_box(slide, '第一章 · 导论', Inches(1), Inches(0.6), Inches(11), Inches(0.5),
             font_size=16, font_color=MUTED, font_name='Consolas')

add_text_box(slide, '什么是交通工程学？', Inches(1), Inches(1.1), Inches(11), Inches(0.8),
             font_size=36, font_color=WHITE, bold=True)

# Left column
add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(2.3), Inches(5.5), Inches(2.3),
          fill_color=DARK_SURFACE, line_color=RGBColor(79, 195, 247))

add_text_box(slide, '学科定义', Inches(1.3), Inches(2.5), Inches(5), Inches(0.5),
             font_size=18, font_color=ACCENT, bold=True)

add_text_box(slide, 
    '交通工程学是研究交通流特性、交通系统\n'
    '规划、设计、管理与控制的综合性工程学科，\n'
    '涉及人、车、路、环境四要素的相互关系\n'
    '与优化配置。',
    Inches(1.3), Inches(3.1), Inches(5), Inches(1.3),
    font_size=15, font_color=TEXT_SOFT, line_spacing=1.3)

add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(4.8), Inches(5.5), Inches(2.3),
          fill_color=DARK_SURFACE, line_color=RGBColor(79, 195, 247))

add_text_box(slide, '研究对象', Inches(1.3), Inches(5), Inches(5), Inches(0.5),
             font_size=18, font_color=ACCENT, bold=True)

bullets = ['交通流特性与规律', '道路通行能力与服务水平', '交通规划与设计方法', '交通管理与控制技术']
add_bullets(slide, bullets, Inches(1.3), Inches(5.4), Inches(5), Inches(1.5), font_size=14, font_color=TEXT_SOFT)

# Right column
add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(2.3), Inches(5.5), Inches(2.3),
          fill_color=DARK_SURFACE, line_color=RGBColor(102, 187, 106))

add_text_box(slide, '发展历程', Inches(7.1), Inches(2.5), Inches(5), Inches(0.5),
             font_size=18, font_color=GREEN, bold=True)

add_text_box(slide,
    '1930s 美国起步 → 1950s 形成体系\n'
    '→ 1970s 引入计算机 → 2000s 智能交通\n'
    '→ 2020s 车路协同与自动驾驶',
    Inches(7.1), Inches(3.1), Inches(5), Inches(1.3),
    font_size=14, font_color=TEXT_SOFT, line_spacing=1.3)

add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(4.8), Inches(5.5), Inches(2.3),
          fill_color=DARK_SURFACE, line_color=RGBColor(255, 183, 77))

add_text_box(slide, '相关学科', Inches(7.1), Inches(5), Inches(5), Inches(0.5),
             font_size=18, font_color=AMBER, bold=True)

add_text_box(slide,
    '城市规划 · 土木工程 · 系统工程\n'
    '运筹学 · 计算机科学\n'
    '环境科学 · 行为科学',
    Inches(7.1), Inches(5.5), Inches(5), Inches(1.3),
    font_size=14, font_color=TEXT_SOFT, line_spacing=1.4)

# ============== SLIDE 5: TRAFFIC FLOW BASICS ==============
slide = prs.slides.add_slide(prs.slide_layouts[6])

bg = add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5),
               fill_color=DARK_BG)

add_text_box(slide, '第二章 · 交通流理论', Inches(1), Inches(0.6), Inches(11), Inches(0.5),
             font_size=16, font_color=MUTED, font_name='Consolas')

add_text_box(slide, '交通流三参数：流量、速度、密度', Inches(1), Inches(1.1), Inches(11), Inches(0.8),
             font_size=32, font_color=WHITE, bold=True)

# Three parameter cards
params = [
    ('Q', '流量 (Flow)', '单位时间内通过某断面的\n车辆数\n\n单位：pcu/h（标准车/小时）\n\nQ = N / T', ACCENT),
    ('V', '速度 (Speed)', '车辆行驶的快慢程度\n\n地点速度 vs 区间速度\n\nV = L / T', GREEN),
    ('K', '密度 (Density)', '单位长度道路上存在的\n车辆数\n\n单位：pcu/km\n\nK = N / L', AMBER),
]

for i, (symbol, title, desc, color) in enumerate(params):
    left = 1 + i * 3.9
    top = 2.3
    
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(3.6), Inches(2.8),
              fill_color=DARK_SURFACE, line_color=RGBColor(79, 195, 247))
    
    # Top accent bar
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(3.6), Inches(0.06),
              fill_color=color)
    
    # Symbol
    add_text_box(slide, symbol, Inches(left + 1.3), Inches(top + 0.2), Inches(1), Inches(0.8),
                 font_size=48, font_color=color, bold=True)
    
    # Title
    add_text_box(slide, title, Inches(left + 0.3), Inches(top + 1), Inches(3), Inches(0.5),
                 font_size=18, font_color=WHITE, bold=True)
    
    # Description
    add_text_box(slide, desc, Inches(left + 0.3), Inches(top + 1.5), Inches(3), Inches(1.2),
                 font_size=13, font_color=TEXT_SOFT, line_spacing=1.3)

# Formula box at bottom
add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(5.5), Inches(11.333), Inches(1.2),
          fill_color=DARK_SURFACE, line_color=RGBColor(79, 195, 247))

add_text_box(slide, '基本关系式：  Q = V × K    —— 交通流理论的核心方程',
             Inches(1), Inches(5.7), Inches(11.333), Inches(0.8),
             font_size=22, font_color=TEXT_SOFT)

add_text_box(slide, 'Q = V × K', Inches(5.2), Inches(5.75), Inches(3), Inches(0.6),
             font_size=28, font_color=ACCENT, font_name='Consolas', bold=True)

# ============== SLIDE 6: FUNDAMENTAL DIAGRAM ==============
slide = prs.slides.add_slide(prs.slide_layouts[6])

bg = add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5),
               fill_color=DARK_BG)

add_text_box(slide, '第二章 · 交通流理论', Inches(1), Inches(0.6), Inches(11), Inches(0.5),
             font_size=16, font_color=MUTED, font_name='Consolas')

add_text_box(slide, '交通流基本图与 Greenshields 模型', Inches(1), Inches(1.1), Inches(11), Inches(0.8),
             font_size=32, font_color=WHITE, bold=True)

# Left column - formulas
add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(2.2), Inches(5.8), Inches(1.4),
          fill_color=DARK_SURFACE, line_color=RGBColor(79, 195, 247))

add_text_box(slide, '速度-密度关系', Inches(1.3), Inches(2.35), Inches(5.3), Inches(0.4),
             font_size=16, font_color=ACCENT, bold=True)
add_text_box(slide, 'V = Vf × (1 - K / Kj)', Inches(1.3), Inches(2.75), Inches(5.3), Inches(0.5),
             font_size=20, font_color=TEXT_SOFT, font_name='Consolas')
add_text_box(slide, 'Vf：自由流速度  |  Kj：阻塞密度', Inches(1.3), Inches(3.2), Inches(5.3), Inches(0.4),
             font_size=13, font_color=MUTED)

add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(3.8), Inches(5.8), Inches(1.4),
          fill_color=DARK_SURFACE, line_color=RGBColor(102, 187, 106))

add_text_box(slide, '流量-密度关系', Inches(1.3), Inches(3.95), Inches(5.3), Inches(0.4),
             font_size=16, font_color=GREEN, bold=True)
add_text_box(slide, 'Q = Vf × K - (Vf / Kj) × K²', Inches(1.3), Inches(4.35), Inches(5.3), Inches(0.5),
             font_size=18, font_color=TEXT_SOFT, font_name='Consolas')
add_text_box(slide, '抛物线关系，存在最大通行能力 Qm', Inches(1.3), Inches(4.8), Inches(5.3), Inches(0.4),
             font_size=13, font_color=MUTED)

add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(5.4), Inches(5.8), Inches(1.6),
          fill_color=DARK_SURFACE, line_color=RGBColor(255, 183, 77))

add_text_box(slide, '关键状态点', Inches(1.3), Inches(5.55), Inches(5.3), Inches(0.4),
             font_size=16, font_color=AMBER, bold=True)

key_points = [
    '自由流状态：K→0, V→Vf, Q→0',
    '最佳密度：Km = Kj / 2，对应最大流量 Qm',
    '阻塞状态：K→Kj, V→0, Q→0',
]
for i, pt in enumerate(key_points):
    add_text_box(slide, f'• {pt}', Inches(1.3), Inches(6 + i * 0.32), Inches(5.3), Inches(0.35),
                 font_size=13, font_color=TEXT_SOFT)

# Right column - diagram placeholder
add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.2), Inches(2.2), Inches(5.2), Inches(4.8),
          fill_color=RGBColor(21, 42, 66), line_color=RGBColor(79, 195, 247))

add_text_box(slide, '交通流基本图示意', Inches(7.5), Inches(2.4), Inches(4.6), Inches(0.4),
             font_size=14, font_color=MUTED)

# Simple diagram
diagram_left = 8.2
diagram_top = 3.2
diagram_width = 3.5
diagram_height = 2.8

# Axes
add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(diagram_left), Inches(diagram_top + diagram_height),
          Inches(diagram_width), Inches(0.03), fill_color=MUTED)
add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(diagram_left), Inches(diagram_top),
          Inches(0.03), Inches(diagram_height), fill_color=MUTED)

# Labels
add_text_box(slide, '密度 K', Inches(diagram_left + diagram_width/2 - 0.4), Inches(diagram_top + diagram_height + 0.15),
             Inches(1), Inches(0.3), font_size=12, font_color=MUTED)
add_text_box(slide, 'Q', Inches(diagram_left - 0.4), Inches(diagram_top + diagram_height/2 - 0.15),
             Inches(0.5), Inches(0.3), font_size=12, font_color=MUTED)

# Parabolic curve (using a series of small rectangles to approximate)
curve_points = [(0, 1), (0.1, 0.88), (0.2, 0.72), (0.3, 0.52), (0.4, 0.32), (0.5, 0.16), 
                (0.6, 0.08), (0.7, 0.16), (0.8, 0.32), (0.9, 0.52), (1, 0.88)]

for i in range(len(curve_points) - 1):
    x1, y1 = curve_points[i]
    x2, y2 = curve_points[i + 1]
    px1 = diagram_left + x1 * diagram_width
    py1 = diagram_top + y1 * diagram_height
    px2 = diagram_left + x2 * diagram_width
    py2 = diagram_top + y2 * diagram_height
    
    line_len = ((px2 - px1)**2 + (py2 - py1)**2)**0.5
    import math
    angle = math.degrees(math.atan2(py2 - py1, px2 - px1))
    
    line = add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(px1), Inches(py1), Inches(line_len), Inches(0.03),
                     fill_color=ACCENT)
    line.rotation = angle

# Max point
add_shape(slide, MSO_SHAPE.OVAL, Inches(diagram_left + diagram_width/2 - 0.08), 
          Inches(diagram_top + 0.08), Inches(0.16), Inches(0.16), fill_color=AMBER)
add_text_box(slide, 'Qm', Inches(diagram_left + diagram_width/2 - 0.2), Inches(diagram_top - 0.2),
             Inches(0.5), Inches(0.3), font_size=11, font_color=AMBER, bold=True)

# Labels on axis
add_text_box(slide, '0', Inches(diagram_left - 0.1), Inches(diagram_top + diagram_height + 0.15),
             Inches(0.3), Inches(0.3), font_size=10, font_color=MUTED)
add_text_box(slide, 'Kj', Inches(diagram_left + diagram_width - 0.1), Inches(diagram_top + diagram_height + 0.15),
             Inches(0.3), Inches(0.3), font_size=10, font_color=MUTED)
add_text_box(slide, 'Km', Inches(diagram_left + diagram_width/2 - 0.15), Inches(diagram_top + diagram_height + 0.15),
             Inches(0.3), Inches(0.3), font_size=10, font_color=ACCENT_SOFT)

add_text_box(slide, 'Q-K 抛物线关系  |  顶点为最大通行能力', Inches(7.5), Inches(6.3), Inches(4.6), Inches(0.4),
             font_size=12, font_color=MUTED)

# ============== SLIDE 7: ROAD CAPACITY ==============
slide = prs.slides.add_slide(prs.slide_layouts[6])

bg = add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5),
               fill_color=DARK_BG)

add_text_box(slide, '第二章 · 交通流理论', Inches(1), Inches(0.6), Inches(11), Inches(0.5),
             font_size=16, font_color=MUTED, font_name='Consolas')

add_text_box(slide, '道路通行能力与服务水平', Inches(1), Inches(1.1), Inches(11), Inches(0.8),
             font_size=32, font_color=WHITE, bold=True)

# Left column
add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(2.2), Inches(5.5), Inches(2.2),
          fill_color=DARK_SURFACE, line_color=RGBColor(79, 195, 247))

add_text_box(slide, '通行能力分类', Inches(1.3), Inches(2.35), Inches(5), Inches(0.5),
             font_size=16, font_color=ACCENT, bold=True)

capacity_items = [
    ('基本通行能力：', '理想条件下的最大流量'),
    ('可能通行能力：', '实际道路条件下的修正值'),
    ('设计通行能力：', '考虑服务水平的折减值'),
]
for i, (label, desc) in enumerate(capacity_items):
    add_text_box(slide, label, Inches(1.3), Inches(2.9 + i * 0.45), Inches(3), Inches(0.4),
                 font_size=13, font_color=WHITE, bold=True)
    add_text_box(slide, desc, Inches(3.2), Inches(2.9 + i * 0.45), Inches(3), Inches(0.4),
                 font_size=13, font_color=TEXT_SOFT)

add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(4.6), Inches(5.5), Inches(1.8),
          fill_color=DARK_SURFACE, line_color=RGBColor(102, 187, 106))

add_text_box(slide, '修正系数', Inches(1.3), Inches(4.75), Inches(5), Inches(0.5),
             font_size=16, font_color=GREEN, bold=True)

add_text_box(slide, 'C = C₀ × fw × fv × fp × ...', Inches(1.3), Inches(5.15), Inches(5), Inches(0.5),
             font_size=16, font_color=TEXT_SOFT, font_name='Consolas')
add_text_box(slide, '车道宽度修正 fw、大型车修正 fv、\n驾驶员特性 fp、坡度修正、视距修正等',
             Inches(1.3), Inches(5.6), Inches(5), Inches(0.7),
             font_size=12, font_color=TEXT_SOFT)

# Right column - LOS
add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(2.2), Inches(5.5), Inches(3),
          fill_color=DARK_SURFACE, line_color=RGBColor(255, 183, 77))

add_text_box(slide, '服务水平 (LOS)', Inches(7.1), Inches(2.35), Inches(5), Inches(0.5),
             font_size=16, font_color=AMBER, bold=True)

los_items = [
    ('A级', '自由流，驾驶自由度大', GREEN),
    ('B级', '稳定流，轻微受限', GREEN),
    ('C级', '稳定流，明显受限', AMBER),
    ('D级', '接近不稳定流', AMBER),
    ('E级', '通行能力极限', RED),
    ('F级', '强制流，严重拥堵', RED),
]

for i, (level, desc, color) in enumerate(los_items):
    top = 2.95 + i * 0.36
    add_text_box(slide, level, Inches(7.1), Inches(top), Inches(0.8), Inches(0.35),
                 font_size=12, font_color=color, bold=True)
    add_text_box(slide, desc, Inches(8), Inches(top), Inches(4), Inches(0.35),
                 font_size=12, font_color=TEXT_SOFT)

add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(5.4), Inches(5.5), Inches(1.5),
          fill_color=DARK_SURFACE, line_color=RGBColor(171, 71, 188))

add_text_box(slide, '典型数值参考', Inches(7.1), Inches(5.55), Inches(5), Inches(0.4),
             font_size=16, font_color=PURPLE, bold=True)

typical_values = [
    '高速公路基本车道：2000-2400 pcu/h',
    '城市主干道：1600-1800 pcu/h',
    '城市次干道：1400-1600 pcu/h',
]
for i, val in enumerate(typical_values):
    add_text_box(slide, val, Inches(7.1), Inches(5.95 + i * 0.3), Inches(5), Inches(0.3),
                 font_size=12, font_color=TEXT_SOFT)

# ============== SLIDE 8: CHAPTER 3 - TRAFFIC SURVEY ==============
slide = prs.slides.add_slide(prs.slide_layouts[6])

bg = add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5),
               fill_color=DARK_BG)

add_text_box(slide, '第三章 · 交通调查与数据分析', Inches(1), Inches(0.6), Inches(11), Inches(0.5),
             font_size=16, font_color=MUTED, font_name='Consolas')

add_text_box(slide, '交通调查类型与方法', Inches(1), Inches(1.1), Inches(11), Inches(0.8),
             font_size=32, font_color=WHITE, bold=True)

survey_types = [
    ('交通量调查', ACCENT,
     '• 连续式调查（24h/7d）\n• 间歇式调查（高峰时段）\n• 自动化采集（线圈/视频）\n• 关键指标：AADT、DHV、K值'),
    ('速度调查', GREEN,
     '• 地点速度测量（雷达/线圈）\n• 区间速度测量（牌照匹配）\n• 行驶速度与行程速度\n• 85%位车速用于限速设定'),
    ('OD 调查', AMBER,
     '• 起讫点调查（Origin-Destination）\n• 家庭访问/路边询问/明信片\n• 手机信令/GPS 大数据\n• OD 表编制与分布预测'),
]

for i, (title, color, content) in enumerate(survey_types):
    left = 1 + i * 3.9
    top = 2.3
    
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(3.6), Inches(4.5),
              fill_color=DARK_SURFACE, line_color=RGBColor(79, 195, 247))
    
    # Top accent bar
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(3.6), Inches(0.06),
              fill_color=color)
    
    add_text_box(slide, title, Inches(left + 0.3), Inches(top + 0.2), Inches(3), Inches(0.5),
                 font_size=18, font_color=color, bold=True)
    add_text_box(slide, content, Inches(left + 0.3), Inches(top + 0.8), Inches(3), Inches(3.5),
                 font_size=13, font_color=TEXT_SOFT, line_spacing=1.4)

# ============== SLIDE 9: CHAPTER 4 - TRAFFIC PLANNING ==============
slide = prs.slides.add_slide(prs.slide_layouts[6])

bg = add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5),
               fill_color=DARK_BG)

add_text_box(slide, '第四章 · 道路交通规划', Inches(1), Inches(0.6), Inches(11), Inches(0.5),
             font_size=16, font_color=MUTED, font_name='Consolas')

add_text_box(slide, '交通规划四阶段法', Inches(1), Inches(1.1), Inches(11), Inches(0.8),
             font_size=32, font_color=WHITE, bold=True)

stages = [
    ('1', '交通生成 (Trip Generation)', 
     '预测各交通小区的出行发生量与吸引量。\n常用方法：回归分析法、交叉分类法、原单位法。\n影响因素：人口、就业、土地利用、汽车保有量等',
     ACCENT),
    ('2', '交通分布 (Trip Distribution)',
     '确定各小区之间的出行交换量。常用模型：\n重力模型（Gravity Model）、Fratar 法。\n核心思想：出行量与产生量/吸引量成正比，\n与阻抗成反比',
     GREEN),
    ('3', '方式划分 (Mode Split)',
     '预测各交通方式的分担率。常用方法：\nLogit 模型、Probit 模型、转移曲线法。\n影响因素：出行时间、费用、舒适度、\n可达性、政策引导',
     AMBER),
    ('4', '交通分配 (Traffic Assignment)',
     '将 OD 出行量分配到路网各路段。常用算法：\n全有全无法、用户均衡 (UE)、系统最优 (SO)。\nWardrop 两大原理是分配的理论基础',
     PURPLE),
]

for i, (num, title, desc, color) in enumerate(stages):
    top = 2.2 + i * 1.2
    
    # Number circle
    circle = add_shape(slide, MSO_SHAPE.OVAL, Inches(1), Inches(top + 0.15), Inches(0.7), Inches(0.7),
                       fill_color=color)
    add_text_box(slide, num, Inches(1.15), Inches(top + 0.22), Inches(0.5), Inches(0.6),
                 font_size=22, font_color=WHITE, bold=True)
    
    # Card background
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.9), Inches(top), Inches(10.5), Inches(1),
              fill_color=DARK_SURFACE, line_color=RGBColor(79, 195, 247))
    
    add_text_box(slide, title, Inches(2.1), Inches(top + 0.05), Inches(10), Inches(0.4),
                 font_size=16, font_color=color, bold=True)
    add_text_box(slide, desc, Inches(2.1), Inches(top + 0.42), Inches(10), Inches(0.5),
                 font_size=12, font_color=TEXT_SOFT, line_spacing=1.2)
    
    # Arrow between stages
    if i < len(stages) - 1:
        add_text_box(slide, '↓', Inches(6.5), Inches(top + 1.05), Inches(0.5), Inches(0.2),
                     font_size=18, font_color=MUTED)

# ============== SLIDE 10: CHAPTER 5 - INTERSECTION DESIGN ==============
slide = prs.slides.add_slide(prs.slide_layouts[6])

bg = add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5),
               fill_color=DARK_BG)

add_text_box(slide, '第五章 · 交叉口设计与信号控制', Inches(1), Inches(0.6), Inches(11), Inches(0.5),
             font_size=16, font_color=MUTED, font_name='Consolas')

add_text_box(slide, '平面交叉口信号控制设计', Inches(1), Inches(1.1), Inches(11), Inches(0.8),
             font_size=32, font_color=WHITE, bold=True)

# Left column
add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(2.2), Inches(5.5), Inches(2.2),
          fill_color=DARK_SURFACE, line_color=RGBColor(79, 195, 247))

add_text_box(slide, '信号配时参数', Inches(1.3), Inches(2.35), Inches(5), Inches(0.5),
             font_size=16, font_color=ACCENT, bold=True)

signal_params = [
    ('周期长度 C：', '信号灯完成一轮变化所需时间'),
    ('绿信比 λ：', '有效绿灯时间与周期之比'),
    ('相位：', '同时获得通行权的交通流组合'),
    ('损失时间 L：', '启动损失 + 清空时间'),
]
for i, (label, desc) in enumerate(signal_params):
    add_text_box(slide, label, Inches(1.3), Inches(2.9 + i * 0.42), Inches(2.8), Inches(0.4),
                 font_size=12, font_color=WHITE, bold=True)
    add_text_box(slide, desc, Inches(3.2), Inches(2.9 + i * 0.42), Inches(3), Inches(0.4),
                 font_size=12, font_color=TEXT_SOFT)

add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(4.6), Inches(5.5), Inches(1.8),
          fill_color=DARK_SURFACE, line_color=RGBColor(102, 187, 106))

add_text_box(slide, 'Webster 配时方法', Inches(1.3), Inches(4.75), Inches(5), Inches(0.5),
             font_size=16, font_color=GREEN, bold=True)
add_text_box(slide, 'C₀ = (1.5L + 5) / (1 - Y)', Inches(1.3), Inches(5.15), Inches(5), Inches(0.5),
             font_size=18, font_color=TEXT_SOFT, font_name='Consolas')
add_text_box(slide, 'L：总损失时间  |  Y：各相位最大流量比之和', Inches(1.3), Inches(5.65), Inches(5), Inches(0.4),
             font_size=12, font_color=MUTED)
add_text_box(slide, '最佳周期使车辆总延误最小', Inches(1.3), Inches(6), Inches(5), Inches(0.3),
             font_size=12, font_color=TEXT_SOFT)

# Right column
add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(2.2), Inches(5.5), Inches(2.2),
          fill_color=DARK_SURFACE, line_color=RGBColor(255, 183, 77))

add_text_box(slide, '交叉口渠化设计', Inches(7.1), Inches(2.35), Inches(5), Inches(0.5),
             font_size=16, font_color=AMBER, bold=True)

channelization = [
    '车道功能划分（直行/左转/右转）',
    '导流岛与标线设计',
    '左转待转区设置',
    '行人与非机动车过街设施',
    '视距三角形保证',
]
for i, item in enumerate(channelization):
    add_text_box(slide, f'• {item}', Inches(7.1), Inches(2.9 + i * 0.36), Inches(5), Inches(0.35),
                 font_size=12, font_color=TEXT_SOFT)

add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(4.6), Inches(5.5), Inches(1.8),
          fill_color=DARK_SURFACE, line_color=RGBColor(239, 83, 80))

add_text_box(slide, '延误分析', Inches(7.1), Inches(4.75), Inches(5), Inches(0.5),
             font_size=16, font_color=RED, bold=True)

delay_items = [
    '控制延误：信号灯引起的停车延误',
    'HCM 延误公式：均匀延误 + 增量延误',
    '服务水平分级：A（<10s）→ F（>80s）',
]
for i, item in enumerate(delay_items):
    add_text_box(slide, f'• {item}', Inches(7.1), Inches(5.2 + i * 0.36), Inches(5), Inches(0.35),
                 font_size=12, font_color=TEXT_SOFT)

# ============== SLIDE 11: CHAPTER 6 - ITS ==============
slide = prs.slides.add_slide(prs.slide_layouts[6])

bg = add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5),
               fill_color=DARK_BG)

add_text_box(slide, '第六章 · 交通管理与智能交通系统', Inches(1), Inches(0.6), Inches(11), Inches(0.5),
             font_size=16, font_color=MUTED, font_name='Consolas')

add_text_box(slide, '智能交通系统 (ITS) 与前沿技术', Inches(1), Inches(1.1), Inches(11), Inches(0.8),
             font_size=32, font_color=WHITE, bold=True)

its_areas = [
    ('🚦', '交通管理系统', ACCENT,
     '• 交通信号自适应控制\n• 区域协调控制 (SCATS/SCOOT)\n• 交通诱导与信息发布\n• 事件检测与应急管理'),
    ('🚗', '车路协同 (V2X)', GREEN,
     '• 车与车通信 (V2V)\n• 车与基础设施通信 (V2I)\n• C-V2X 与 5G 通信\n• 自动驾驶协同感知\n• 智能网联汽车测试'),
    ('📱', 'MaaS 与大数据', AMBER,
     '• 出行即服务 (Mobility as a Service)\n• 手机信令/公交 IC 卡数据\n• 浮动车 GPS 轨迹分析\n• 人工智能在交通中的应用'),
]

for i, (emoji, title, color, content) in enumerate(its_areas):
    left = 1 + i * 3.9
    top = 2.3
    
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(3.6), Inches(4.5),
              fill_color=DARK_SURFACE, line_color=RGBColor(79, 195, 247))
    
    # Top accent bar
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(3.6), Inches(0.06),
              fill_color=color)
    
    add_text_box(slide, emoji, Inches(left + 0.2), Inches(top + 0.15), Inches(0.8), Inches(0.5),
                 font_size=24)
    add_text_box(slide, title, Inches(left + 1), Inches(top + 0.15), Inches(2.4), Inches(0.5),
                 font_size=16, font_color=color, bold=True)
    add_text_box(slide, content, Inches(left + 0.3), Inches(top + 0.7), Inches(3), Inches(3.5),
                 font_size=13, font_color=TEXT_SOFT, line_spacing=1.3)

# ============== SLIDE 12: SUMMARY ==============
slide = prs.slides.add_slide(prs.slide_layouts[6])

bg = add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5),
               fill_color=DARK_BG)

add_text_box(slide, '课程总结', Inches(1), Inches(0.6), Inches(11), Inches(0.5),
             font_size=16, font_color=MUTED, font_name='Consolas')

add_text_box(slide, '核心知识体系回顾', Inches(1), Inches(1.1), Inches(11), Inches(0.8),
             font_size=40, font_color=WHITE, bold=True)

summary_items = [
    ('', '交通流理论', '掌握流量-速度-密度关系，\n理解通行能力与服务水平'),
    ('📋', '交通调查', '熟练运用各类交通调查方法，\n掌握数据处理与分析技能'),
    ('🗺️', '交通规划', '理解四阶段法原理，能够进行\n交通需求预测与方案评价'),
    ('🚦', '信号控制', '掌握交叉口渠化设计与\n信号配时优化方法'),
]

for i, (emoji, title, desc) in enumerate(summary_items):
    left = 1 + (i % 2) * 5.8
    top = 2.3 + (i // 2) * 2
    
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(5.5), Inches(1.8),
              fill_color=DARK_SURFACE, line_color=RGBColor(79, 195, 247))
    
    add_text_box(slide, emoji, Inches(left + 0.3), Inches(top + 0.1), Inches(0.8), Inches(0.5),
                 font_size=24)
    add_text_box(slide, title, Inches(left + 1.2), Inches(top + 0.1), Inches(4), Inches(0.5),
                 font_size=18, font_color=WHITE, bold=True)
    add_text_box(slide, desc, Inches(left + 0.3), Inches(top + 0.7), Inches(5), Inches(1),
                 font_size=14, font_color=TEXT_SOFT)

# Bottom
add_text_box(slide, '理论 + 实践 + 创新 = 现代交通工程人才', Inches(1), Inches(6.5), Inches(6), Inches(0.5),
             font_size=16, font_color=MUTED, font_name='Consolas')
add_text_box(slide, '交通工程学', Inches(11), Inches(6.5), Inches(2), Inches(0.5),
             font_size=16, font_color=MUTED)

# ============== SLIDE 13: HOMEWORK ==============
slide = prs.slides.add_slide(prs.slide_layouts[6])

bg = add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5),
               fill_color=DARK_BG)

add_text_box(slide, '课后任务', Inches(1), Inches(0.6), Inches(11), Inches(0.5),
             font_size=16, font_color=MUTED, font_name='Consolas')

add_text_box(slide, '思考题与实践作业', Inches(1), Inches(1.1), Inches(11), Inches(0.8),
             font_size=36, font_color=WHITE, bold=True)

# Thinking questions
add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(2.2), Inches(5.5), Inches(3),
          fill_color=DARK_SURFACE, line_color=RGBColor(79, 195, 247))

add_text_box(slide, '📝 思考题', Inches(1.3), Inches(2.35), Inches(5), Inches(0.5),
             font_size=16, font_color=ACCENT, bold=True)

thinking_qs = [
    '1. 交通流三参数之间有什么关系？画出基本图并解释',
    '2. 通行能力与服务水平有何区别与联系？',
    '3. 四阶段法中每个阶段的作用是什么？',
    '4. 如何优化一个拥堵交叉口的信号配时？',
    '5. 智能交通系统如何解决城市交通问题？',
]
for i, q in enumerate(thinking_qs):
    add_text_box(slide, q, Inches(1.3), Inches(2.9 + i * 0.42), Inches(5), Inches(0.4),
                 font_size=12, font_color=TEXT_SOFT)

# Practice assignments
add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(2.2), Inches(5.5), Inches(3),
          fill_color=DARK_SURFACE, line_color=RGBColor(102, 187, 106))

add_text_box(slide, ' 实践作业', Inches(7.1), Inches(2.35), Inches(5), Inches(0.5),
             font_size=16, font_color=GREEN, bold=True)

practice_items = [
    '1. 选择一个交叉口，进行 15 分钟交通量调查',
    '2. 使用 Webster 方法计算最佳信号周期',
    '3. 分析调查路段的服务水平等级',
    '4. 提出改善建议并撰写调查报告',
    '5. 查阅一篇智能交通相关论文并做课堂分享',
]
for i, item in enumerate(practice_items):
    add_text_box(slide, item, Inches(7.1), Inches(2.9 + i * 0.42), Inches(5), Inches(0.4),
                 font_size=12, font_color=TEXT_SOFT)

# Recommended reading
add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(5.4), Inches(11.333), Inches(1.5),
          fill_color=DARK_SURFACE, line_color=RGBColor(255, 183, 77))

add_text_box(slide, '📚 推荐阅读', Inches(1.3), Inches(5.55), Inches(11), Inches(0.4),
             font_size=16, font_color=AMBER, bold=True)

readings = [
    '《交通工程学》（第 5 版），任福田主编，人民交通出版社',
    '《交通流理论》，王炜等，科学出版社',
    'Highway Capacity Manual (HCM 6th Edition), TRB',
]
for i, r in enumerate(readings):
    add_text_box(slide, f'• {r}', Inches(1.3), Inches(6 + i * 0.28), Inches(11), Inches(0.3),
                 font_size=12, font_color=TEXT_SOFT)

# ============== SLIDE 14: AI MARKER / THANK YOU ==============
slide = prs.slides.add_slide(prs.slide_layouts[6])

bg = add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5),
               fill_color=DARK_BG)

# Decorative elements
for i in range(8):
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(11 + i * 0.15), Inches(0), Inches(0.04), Inches(7.5),
              fill_color=RGBColor(79, 195, 247))

add_text_box(slide, '感谢观看', Inches(1), Inches(2.2), Inches(11), Inches(1.2),
             font_size=56, font_color=WHITE, bold=True)

add_text_box(slide, '交通工程学 —— 让出行更美好，让城市更智慧', Inches(1), Inches(3.6), Inches(10), Inches(0.8),
             font_size=24, font_color=TEXT_SOFT)

# AI marker
add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4), Inches(5.2), Inches(5.333), Inches(0.8),
          fill_color=DARK_SURFACE, line_color=RGBColor(79, 195, 247))

add_text_box(slide, '本文由 AI 辅助创作 · 交通工程学课程 PPT · 2026 年',
             Inches(4), Inches(5.35), Inches(5.333), Inches(0.5),
             font_size=14, font_color=MUTED)

# ============== SAVE ==============
output_path = '/home/admin/.openclaw/workspace/交通工程学_授课PPT.pptx'
prs.save(output_path)
print(f'✅ PPT generated: {output_path}')
print(f'📊 Total slides: {len(prs.slides)}')

# Print slide info
for i, slide in enumerate(prs.slides):
    print(f'  Slide {i+1}: {len(slide.shapes)} shapes')
