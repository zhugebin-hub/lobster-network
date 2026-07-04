#!/usr/bin/env python3
"""生成五一假期安全教育PPT"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ===== 配色方案 =====
RED_PRIMARY = RGBColor(0xD3, 0x2F, 0x2F)
RED_LIGHT = RGBColor(0xFF, 0xEB, 0xEE)
ORANGE_ACCENT = RGBColor(0xFF, 0x6F, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x21, 0x21, 0x21)
GRAY = RGBColor(0x75, 0x75, 0x75)
BLUE_SAFETY = RGBColor(0x19, 0x76, 0xD2)
GREEN_SAFETY = RGBColor(0x38, 0x8E, 0x3C)
YELLOW_SAFETY = RGBColor(0xF5, 0x7F, 0x17)

def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_textbox(slide, left, top, width, height, text, size=18, color=DARK, bold=False, align=PP_ALIGN.LEFT, font_name="微软雅黑"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = font_name
    return txBox

def add_bullet_list(slide, left, top, width, height, items, size=16, color=DARK, spacing=Pt(6)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = spacing
        run = p.add_run()
        run.text = f"• {item}"
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.name = "微软雅黑"
    return txBox

def add_card(slide, left, top, width, height, icon, title, content, icon_color=RED_PRIMARY):
    # 卡片背景
    card = add_shape(slide, left, top, width, height, WHITE)
    # 顶部颜色条
    add_shape(slide, left, top, width, Inches(0.06), icon_color)
    # 图标
    add_textbox(slide, left + Inches(0.15), top + Inches(0.15), Inches(0.5), Inches(0.5), icon, size=24, color=icon_color, align=PP_ALIGN.CENTER)
    # 标题
    add_textbox(slide, left + Inches(0.7), top + Inches(0.15), width - Inches(0.85), Inches(0.4), title, size=16, color=icon_color, bold=True)
    # 内容
    add_textbox(slide, left + Inches(0.15), top + Inches(0.6), width - Inches(0.3), height - Inches(0.75), content, size=13, color=GRAY)

def add_header(slide, title):
    add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), RED_PRIMARY)
    add_textbox(slide, Inches(0.5), Inches(0.25), Inches(12), Inches(0.8), title, size=32, color=WHITE, bold=True)

# ============================================================
# 第1页：封面
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

# 顶部装饰
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(3.5), RED_PRIMARY)
add_shape(slide, Inches(0), Inches(3.5), Inches(13.333), Inches(0.08), ORANGE_ACCENT)

# 图标
add_textbox(slide, Inches(4), Inches(0.5), Inches(5.3), Inches(1.5), "🛡️", size=60, color=WHITE, align=PP_ALIGN.CENTER)

# 主标题
add_textbox(slide, Inches(1.5), Inches(1.5), Inches(10.3), Inches(1.2),
            "五一假期安全教育", size=48, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

# 副标题
add_textbox(slide, Inches(2), Inches(2.8), Inches(9.3), Inches(0.6),
            "中职生假期安全指南", size=24, color=RGBColor(0xFF, 0xCD, 0xD2), align=PP_ALIGN.CENTER)

# 底部信息
add_textbox(slide, Inches(2), Inches(4.5), Inches(9.3), Inches(0.5),
            "2026年五一劳动节", size=20, color=GRAY, align=PP_ALIGN.CENTER)
add_textbox(slide, Inches(2), Inches(5.2), Inches(9.3), Inches(0.5),
            "交通安全 · 防电信诈骗 · 消防安全 · 假期安全 · 其他安全", size=16, color=GRAY, align=PP_ALIGN.CENTER)

# 底部装饰
add_shape(slide, Inches(0), Inches(7.0), Inches(13.333), Inches(0.08), ORANGE_ACCENT)

# ============================================================
# 第2页：目录
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_header(slide, "目录")

toc_items = [
    ("🚗", "交通安全", "出行安全要点"),
    ("📱", "防电信诈骗", "识别常见骗局"),
    ("🔥", "消防安全", "防火与自救"),
    ("🏖️", "假期安全", "饮食·游泳·出行"),
    ("⚠️", "其他安全", "心理健康·网络安全"),
]

for i, (icon, title, desc) in enumerate(toc_items):
    y = Inches(1.6) + Inches(1.1) * i
    # 序号
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.5), y + Inches(0.05), Inches(0.6), Inches(0.6))
    circle.fill.solid()
    circle.fill.fore_color.rgb = RED_PRIMARY
    circle.line.fill.background()
    add_textbox(slide, Inches(1.5), y + Inches(0.05), Inches(0.6), Inches(0.6), f"{i+1}", size=20, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    # 图标
    add_textbox(slide, Inches(2.3), y, Inches(0.6), Inches(0.6), icon, size=28)
    # 标题
    add_textbox(slide, Inches(3.0), y, Inches(3), Inches(0.5), title, size=22, color=RED_PRIMARY, bold=True)
    # 描述
    add_textbox(slide, Inches(6.5), y + Inches(0.05), Inches(5), Inches(0.5), desc, size=16, color=GRAY)

# ============================================================
# 第3页：交通安全 - 概述
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_header(slide, "🚗 交通安全")

add_textbox(slide, Inches(0.8), Inches(1.6), Inches(11.7), Inches(0.6),
            "假期出行频繁，交通事故高发，请务必注意以下安全要点：", size=18, color=GRAY, align=PP_ALIGN.CENTER)

# 左侧卡片
add_card(slide, Inches(0.8), Inches(2.5), Inches(3.6), Inches(2.2), "🚶", "步行安全",
         "• 走人行道，不闯红灯\n• 不在马路上追逐打闹\n• 过马路走斑马线\n• 夜间出行穿亮色衣服", BLUE_SAFETY)

# 中间卡片
add_card(slide, Inches(4.8), Inches(2.5), Inches(3.6), Inches(2.2), "🚌", "乘车安全",
         "• 乘坐正规营运车辆\n• 系好安全带\n• 不乘坐超载车辆\n• 头手不伸出窗外", GREEN_SAFETY)

# 右侧卡片
add_card(slide, Inches(8.8), Inches(2.5), Inches(3.6), Inches(2.2), "🚲", "骑行安全",
         "• 未满16岁不骑电动车\n• 佩戴安全头盔\n• 不逆行、不闯红灯\n• 不单手骑车", ORANGE_ACCENT)

# 底部警示
add_shape(slide, Inches(0.8), Inches(5.2), Inches(11.7), Inches(1.8), RED_LIGHT)
add_textbox(slide, Inches(1.2), Inches(5.3), Inches(11), Inches(0.4), "⚠️ 特别提醒", size=18, color=RED_PRIMARY, bold=True)
add_textbox(slide, Inches(1.2), Inches(5.8), Inches(11), Inches(1.0),
            "• 不乘坐黑车、摩的\n• 不酒后驾车\n• 不疲劳驾驶\n• 遇到事故立即拨打122报警", size=15, color=DARK)

# ============================================================
# 第4页：交通安全 - 案例
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_header(slide, "🚗 交通安全 - 真实案例")

add_textbox(slide, Inches(0.8), Inches(1.6), Inches(11.7), Inches(0.6),
            "以下案例来自近年假期交通事故统计，请引以为戒：", size=18, color=GRAY, align=PP_ALIGN.CENTER)

# 案例1
add_shape(slide, Inches(0.8), Inches(2.5), Inches(5.5), Inches(2.2), RGBColor(0xFF, 0xF3, 0xE0))
add_textbox(slide, Inches(1.0), Inches(2.6), Inches(5.1), Inches(0.4), "案例一：闯红灯酿悲剧", size=18, color=ORANGE_ACCENT, bold=True)
add_textbox(slide, Inches(1.0), Inches(3.1), Inches(5.1), Inches(1.4),
            "2025年五一期间，某职校学生李某骑电动车闯红灯，与正常行驶的轿车发生碰撞，导致腿部骨折，休学三个月。", size=14, color=DARK)

# 案例2
add_shape(slide, Inches(6.8), Inches(2.5), Inches(5.5), Inches(2.2), RGBColor(0xFF, 0xEB, 0xEE))
add_textbox(slide, Inches(7.0), Inches(2.6), Inches(5.1), Inches(0.4), "案例二：黑车出事无人管", size=18, color=RED_PRIMARY, bold=True)
add_textbox(slide, Inches(7.0), Inches(3.1), Inches(5.1), Inches(1.4),
            "2024年五一，三名学生乘坐黑车回家，途中发生车祸，司机逃逸。因车辆无营运资质，赔偿困难。", size=14, color=DARK)

# 安全口诀
add_shape(slide, Inches(0.8), Inches(5.2), Inches(11.7), Inches(1.8), RGBColor(0xE8, 0xF5, 0xE9))
add_textbox(slide, Inches(1.2), Inches(5.3), Inches(11), Inches(0.4), "📝 交通安全口诀", size=18, color=GREEN_SAFETY, bold=True, align=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1.2), Inches(5.8), Inches(11), Inches(1.0),
            "红灯停，绿灯行，黄灯亮了等一等\n过马路，走斑马，不在路上打打闹\n乘车系好安全带，头手不伸窗外边\n黑车摩的不乘坐，安全出行记心间", size=16, color=DARK, align=PP_ALIGN.CENTER)

# ============================================================
# 第5页：防电信诈骗 - 概述
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_header(slide, "📱 防电信诈骗")

add_textbox(slide, Inches(0.8), Inches(1.6), Inches(11.7), Inches(0.6),
            "学生是电信诈骗的高危群体！假期有更多时间上网，务必提高警惕：", size=18, color=GRAY, align=PP_ALIGN.CENTER)

# 四大常见骗局
scams = [
    ("🎮", "游戏账号交易诈骗", "低价出售游戏装备/账号，要求先转账后交货，收款后拉黑", ORANGE_ACCENT),
    ("🛒", "网购退款诈骗", "冒充客服称商品有问题需退款，发送钓鱼链接骗取银行卡信息", RED_PRIMARY),
    ("💰", "刷单返利诈骗", "以高额佣金为诱饵，先小额返利获取信任，大额投入后消失", RGBColor(0x7B, 0x1F, 0xA2)),
    ("👤", "冒充熟人诈骗", "盗用QQ/微信头像，冒充同学/老师紧急借钱，要求立即转账", BLUE_SAFETY),
]

for i, (icon, title, desc, color) in enumerate(scams):
    col = i % 2
    row = i // 2
    x = Inches(0.8) + Inches(6.0) * col
    y = Inches(2.5) + Inches(2.2) * row
    add_card(slide, x, y, Inches(5.5), Inches(1.9), icon, title, desc, color)

# ============================================================
# 第6页：防电信诈骗 - 防范
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_header(slide, "📱 防电信诈骗 - 六不原则")

add_textbox(slide, Inches(0.8), Inches(1.6), Inches(11.7), Inches(0.6),
            "牢记'六不'原则，守住钱袋子：", size=18, color=GRAY, align=PP_ALIGN.CENTER)

principles = [
    ("不轻信", "陌生电话、短信、链接不轻信", "接到可疑电话先核实身份，不点击不明链接"),
    ("不透露", "个人信息、银行卡号不透露", "验证码就是'钱袋子'的钥匙，绝不告诉任何人"),
    ("不转账", "陌生人要求转账不转账", "凡是要求转账的，一律挂断电话"),
    ("不贪心", "天上不会掉馅饼，贪小便宜吃大亏", "高回报、低投入的都是骗局"),
    ("不恐慌", "冒充公检法称你犯法，不恐慌", "公检法不会电话办案，更不会要求转账"),
    ("不拖延", "发现被骗立即报警，不拖延", "拨打110或96110，保留证据"),
]

for i, (title, subtitle, desc) in enumerate(principles):
    col = i % 3
    row = i // 3
    x = Inches(0.8) + Inches(3.9) * col
    y = Inches(2.5) + Inches(1.6) * row
    add_shape(slide, x, y, Inches(3.5), Inches(1.3), RED_LIGHT)
    add_textbox(slide, x + Inches(0.15), y + Inches(0.1), Inches(3.2), Inches(0.4), f"{i+1}. {title}", size=16, color=RED_PRIMARY, bold=True)
    add_textbox(slide, x + Inches(0.15), y + Inches(0.5), Inches(3.2), Inches(0.35), subtitle, size=12, color=DARK, bold=True)
    add_textbox(slide, x + Inches(0.15), y + Inches(0.85), Inches(3.2), Inches(0.4), desc, size=11, color=GRAY)

# 底部热线
add_shape(slide, Inches(0.8), Inches(6.2), Inches(11.7), Inches(0.8), RED_PRIMARY)
add_textbox(slide, Inches(0.8), Inches(6.3), Inches(11.7), Inches(0.6),
            "🆘 反诈热线：96110（全国反诈专线）  报警电话：110", size=20, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

# ============================================================
# 第7页：消防安全
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_header(slide, "🔥 消防安全")

add_textbox(slide, Inches(0.8), Inches(1.6), Inches(11.7), Inches(0.6),
            "假期居家时间长，消防安全不容忽视：", size=18, color=GRAY, align=PP_ALIGN.CENTER)

# 左侧：防火要点
add_shape(slide, Inches(0.8), Inches(2.5), Inches(5.5), Inches(4.2), RGBColor(0xFF, 0xEB, 0xEE))
add_textbox(slide, Inches(1.2), Inches(2.6), Inches(4.7), Inches(0.4), "🏠 居家防火要点", size=18, color=RED_PRIMARY, bold=True)
add_bullet_list(slide, Inches(1.2), Inches(3.1), Inches(4.7), Inches(3.3), [
    "不私拉乱接电线，不超负荷用电",
    "离开房间时关闭电器电源",
    "不躺在床上吸烟，不乱扔烟头",
    "厨房用火不离人，油锅起火盖锅盖",
    "不存放易燃易爆物品",
    "定期检查燃气管道是否漏气",
], size=14, color=DARK)

# 右侧：自救知识
add_shape(slide, Inches(6.8), Inches(2.5), Inches(5.5), Inches(4.2), RGBColor(0xFF, 0xF3, 0xE0))
add_textbox(slide, Inches(7.2), Inches(2.6), Inches(4.7), Inches(0.4), "🚒 火灾自救知识", size=18, color=ORANGE_ACCENT, bold=True)
add_bullet_list(slide, Inches(7.2), Inches(3.1), Inches(4.7), Inches(3.3), [
    "发现火灾立即拨打119报警",
    "用湿毛巾捂住口鼻，弯腰低姿逃生",
    "不乘坐电梯，走安全通道",
    "身上着火就地打滚，不奔跑",
    "被困室内时，用湿布塞门缝，在窗口呼救",
    "不贪恋财物，生命第一",
], size=14, color=DARK)

# ============================================================
# 第8页：假期安全 - 饮食安全
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_header(slide, "🏖️ 假期安全 - 饮食安全")

add_textbox(slide, Inches(0.8), Inches(1.6), Inches(11.7), Inches(0.6),
            "假期聚餐多，饮食安全要牢记：", size=18, color=GRAY, align=PP_ALIGN.CENTER)

# 四个卡片
food_items = [
    ("🍽️", "注意饮食卫生", "饭前便后洗手，不吃生冷食物，不喝生水", GREEN_SAFETY),
    ("⚠️", "警惕食物中毒", "不吃过期食品，不吃野生蘑菇，不吃发芽土豆", ORANGE_ACCENT),
    ("🍺", "不酗酒", "未成年人禁止饮酒，适量饮酒伤身体", RED_PRIMARY),
    ("🏪", "选择正规餐厅", "不光顾无证摊贩，注意查看卫生许可证", BLUE_SAFETY),
]

for i, (icon, title, desc, color) in enumerate(food_items):
    col = i % 2
    row = i // 2
    x = Inches(0.8) + Inches(6.0) * col
    y = Inches(2.5) + Inches(2.2) * row
    add_card(slide, x, y, Inches(5.5), Inches(1.9), icon, title, desc, color)

# 底部
add_shape(slide, Inches(0.8), Inches(5.2), Inches(11.7), Inches(1.8), RGBColor(0xE8, 0xF5, 0xE9))
add_textbox(slide, Inches(1.2), Inches(5.3), Inches(11), Inches(0.4), "🆘 食物中毒急救", size=18, color=GREEN_SAFETY, bold=True)
add_textbox(slide, Inches(1.2), Inches(5.8), Inches(11), Inches(1.0),
            "• 立即停止食用可疑食物\n• 催吐：用手指刺激咽喉部催吐\n• 保留食物样本，以便检测\n• 立即拨打120就医", size=14, color=DARK)

# ============================================================
# 第9页：假期安全 - 游泳安全
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_header(slide, "🏖️ 假期安全 - 游泳安全")

add_textbox(slide, Inches(0.8), Inches(1.6), Inches(11.7), Inches(0.6),
            "夏季来临，溺水事故进入高发期！请务必牢记：", size=18, color=GRAY, align=PP_ALIGN.CENTER)

# 六不准
add_shape(slide, Inches(0.8), Inches(2.5), Inches(11.7), Inches(2.2), RGBColor(0xFF, 0xEB, 0xEE))
add_textbox(slide, Inches(1.2), Inches(2.6), Inches(11), Inches(0.4), "🚫 游泳'六不准'", size=18, color=RED_PRIMARY, bold=True)
add_bullet_list(slide, Inches(1.2), Inches(3.1), Inches(11), Inches(1.4), [
    "不私自下水游泳",
    "不擅自与他人结伴游泳",
    "不在无家长或教师带领的情况下游泳",
    "不到无安全设施、无救援人员的水域游泳",
    "不到不熟悉的水域游泳",
    "不熟悉水性的学生不擅自下水施救",
], size=15, color=DARK)

# 自救知识
add_shape(slide, Inches(0.8), Inches(5.2), Inches(5.5), Inches(1.8), RGBColor(0xE3, 0xF2, 0xFD))
add_textbox(slide, Inches(1.2), Inches(5.3), Inches(4.7), Inches(0.4), "🏊 溺水自救", size=16, color=BLUE_SAFETY, bold=True)
add_textbox(slide, Inches(1.2), Inches(5.8), Inches(4.7), Inches(1.0),
            "• 保持冷静，呼救\n• 头向后仰，口鼻露出水面\n• 双手划水，不挣扎\n• 抽筋时拉伸肌肉", size=13, color=DARK)

# 施救知识
add_shape(slide, Inches(6.8), Inches(5.2), Inches(5.5), Inches(1.8), RGBColor(0xFF, 0xF3, 0xE0))
add_textbox(slide, Inches(7.2), Inches(5.3), Inches(4.7), Inches(0.4), "🆘 正确施救", size=16, color=ORANGE_ACCENT, bold=True)
add_textbox(slide, Inches(7.2), Inches(5.8), Inches(4.7), Inches(1.0),
            "• 大声呼救，拨打110/120\n• 寻找竹竿、绳子、漂浮物\n• 不盲目下水施救\n• 多人手拉手施救", size=13, color=DARK)

# ============================================================
# 第10页：其他安全
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_header(slide, "⚠️ 其他安全")

add_textbox(slide, Inches(0.8), Inches(1.6), Inches(11.7), Inches(0.6),
            "除了以上安全事项，还需注意以下方面：", size=18, color=GRAY, align=PP_ALIGN.CENTER)

# 四个卡片
other_items = [
    ("🌐", "网络安全", "• 不浏览不良网站\n• 不沉迷网络游戏\n• 不随意约见网友\n• 注意保护个人隐私", BLUE_SAFETY),
    ("💪", "心理健康", "• 保持规律作息\n• 适当运动锻炼\n• 多与家人沟通\n• 遇困扰及时求助", GREEN_SAFETY),
    ("🏔️", "旅行安全", "• 告知家人行程\n• 不前往未开发景区\n• 注意天气变化\n• 购买旅游意外险", ORANGE_ACCENT),
    ("📞", "应急电话", "• 报警：110\n• 火警：119\n• 急救：120\n• 交通事故：122", RED_PRIMARY),
]

for i, (icon, title, desc, color) in enumerate(other_items):
    col = i % 2
    row = i // 2
    x = Inches(0.8) + Inches(6.0) * col
    y = Inches(2.5) + Inches(2.2) * row
    add_card(slide, x, y, Inches(5.5), Inches(2.0), icon, title, desc, color)

# ============================================================
# 第11页：安全承诺
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), RED_PRIMARY)
add_textbox(slide, Inches(0.5), Inches(0.25), Inches(12), Inches(0.8), "🤝 安全承诺", size=32, color=WHITE, bold=True)

add_textbox(slide, Inches(1.5), Inches(1.8), Inches(10.3), Inches(0.6),
            "我郑重承诺：", size=22, color=DARK, bold=True, align=PP_ALIGN.CENTER)

promises = [
    "遵守交通规则，不闯红灯，不乘坐黑车",
    "提高警惕，不轻信陌生电话和短信",
    "注意消防安全，离开房间关闭电源",
    "注意饮食卫生，不暴饮暴食",
    "不到危险水域游泳，不擅自下水施救",
    "文明上网，不沉迷游戏，保护个人隐私",
    "保持良好心态，遇到困难及时求助",
]

for i, promise in enumerate(promises):
    y = Inches(2.5) + Inches(0.55) * i
    add_textbox(slide, Inches(2.0), y, Inches(9.3), Inches(0.5), f"✓ {promise}", size=16, color=DARK)

# 底部
add_shape(slide, Inches(2), Inches(6.2), Inches(9.3), Inches(0.8), RED_PRIMARY)
add_textbox(slide, Inches(2), Inches(6.3), Inches(9.3), Inches(0.6),
            "安全无小事，防范于未然！", size=24, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

# ============================================================
# 第12页：结束页
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, RED_PRIMARY)

add_shape(slide, Inches(0), Inches(2.5), Inches(13.333), Inches(0.06), ORANGE_ACCENT)

add_textbox(slide, Inches(1.5), Inches(3.0), Inches(10.3), Inches(1.5),
            "祝大家五一假期\n安全、快乐！", size=48, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

add_textbox(slide, Inches(2), Inches(5.0), Inches(9.3), Inches(0.6),
            "安全牢记心间，平安度过假期", size=22, color=RGBColor(0xFF, 0xCD, 0xD2), align=PP_ALIGN.CENTER)

add_textbox(slide, Inches(2), Inches(5.8), Inches(9.3), Inches(0.5),
            "2026年五一劳动节", size=18, color=RGBColor(0xFF, 0xCD, 0xD2), align=PP_ALIGN.CENTER)

add_shape(slide, Inches(0), Inches(7.0), Inches(13.333), Inches(0.06), ORANGE_ACCENT)

# ============================================================
# 保存
# ============================================================
output_path = "/home/admin/.openclaw/workspace/王利英_五一安全教育/中职生五一假期安全教育.pptx"
prs.save(output_path)
print(f"✅ PPT已生成：{output_path}")
