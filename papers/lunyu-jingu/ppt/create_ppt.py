#!/usr/bin/env python3
"""《论语今文经学与古文经学比较研究》PPT 生成 + 讲稿"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree
from datetime import date

# 作者信息
AUTHOR = '冯诚禾'
COURSE = 'AI赋能'
ADVISOR = '诸葛斌'
TITLE_TEXT = '论语今文经学与古文经学比较研究'

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BG    = RGBColor(0x1B, 0x2A, 0x4A)
BG2   = RGBColor(0x2C, 0x3E, 0x6B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GOLD  = RGBColor(0xE8, 0xC5, 0x47)
GOLD2 = RGBColor(0xC4, 0xA0, 0x2A)
LBLUE = RGBColor(0xA8, 0xD8, 0xEA)
WARM  = RGBColor(0xF5, 0xF0, 0xE8)
GREEN = RGBColor(0x90, 0xEE, 0x90)
RED   = RGBColor(0xFF, 0x99, 0x99)
GREY  = RGBColor(0x99, 0x99, 0x99)

def _run(paragraph, text, name='微软雅黑', size=18, bold=False, color=None, italic=False):
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.bold = bold
    run.italic = italic
    if color: run.font.color.rgb = color
    run.font.name = name
    # 设置中文字体
    rPr = run.font._element
    rFonts = etree.SubElement(rPr, qn('a:rFonts'))
    rFonts.set(qn('a:eastAsia'), name)
    rFonts.set(qn('a:ascii'), name)
    return run

def fill_bg(slide, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

def rect(slide, l, t, w, h, color):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()
    return s

def gold_bar(slide, l, t, w, h=Inches(0.04)):
    rect(slide, l, t, w, h, GOLD)

def _tb(slide, l, t, w, h, text, name='微软雅黑', size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tb.text_frame.word_wrap = True
    p = tb.text_frame.paragraphs[0]
    p.alignment = align
    _run(p, text, name, size, bold, color)
    return tb

def _mb(slide, l, t, w, h, lines):
    """lines: (text, size, bold, color, align)"""
    tb = slide.shapes.add_textbox(l, t, w, h)
    tb.text_frame.word_wrap = True
    for i, item in enumerate(lines):
        text = item[0]
        sz = item[1] if len(item)>1 else 18
        bd = item[2] if len(item)>2 else False
        cl = item[3] if len(item)>3 else WHITE
        al = item[4] if len(item)>4 else PP_ALIGN.LEFT
        p = tb.text_frame.paragraphs[0] if i==0 else tb.text_frame.add_paragraph()
        p.alignment = al
        p.space_after = Pt(3); p.space_before = Pt(2)
        _run(p, text, '微软雅黑', sz, bd, cl)
    return tb

def title_bar(slide, text):
    rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), BG)
    gold_bar(slide, Inches(0), Inches(1.2), Inches(13.333))
    _tb(slide, Inches(0.8), Inches(0.25), Inches(11.5), Inches(0.8),
        text, name='黑体', size=26, bold=True, color=GOLD)

def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text

def blank():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(s, BG)
    return s

# ================================================================
# SLIDE 0: 封面
# ================================================================
s = blank()
rect(s, Inches(0), Inches(2.8), Inches(13.333), Inches(0.06), GOLD)
rect(s, Inches(0), Inches(5.3), Inches(13.333), Inches(0.04), GOLD2)
_tb(s, Inches(1), Inches(1.0), Inches(11.3), Inches(1.6),
    TITLE_TEXT, name='黑体', size=36, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
_mb(s, Inches(1.5), Inches(3.2), Inches(10.3), Inches(2.0), [
    ('从版本源流到诠释范式的千年之争', 22, False, LBLUE, PP_ALIGN.CENTER),
    ('', 6),
    (f'作者：{AUTHOR}    课程：{COURSE}    导师：{ADVISOR}', 17, False, WARM, PP_ALIGN.CENTER),
    ('', 6),
    (date.today().strftime('%Y年%-m月%-d日'), 14, False, GREY, PP_ALIGN.CENTER),
])
notes(s, f"""各位老师、同学，大家好。

我是{AUTHOR}，本次汇报的课程是《{COURSE}》，在{ADVISOR}老师的指导下完成。

今天我汇报的题目是《{TITLE_TEXT}》。

《论语》作为儒家经典的核心文献，在汉代经历了两条截然不同的传承路径——今文经学与古文经学。这两条路径不仅仅是文本版本的差异，更代表了两种完全不同的经典诠释范式。

在接下来约15分钟的汇报中，我将从版本源流、学术特征、诠释方法、个案分析和历史演变五个方面，系统地比较今文经学与古文经学的异同，并探讨这一学术争辩的当代启示。

我们先从《论语》在汉代的版本分化说起。""")

# ================================================================
# SLIDE 1: 目录
# ================================================================
s = blank()
_tb(s, Inches(0.8), Inches(0.3), Inches(5), Inches(0.8),
    '汇 报 提 纲', name='黑体', size=28, bold=True, color=GOLD)
gold_bar(s, Inches(0.8), Inches(1.1), Inches(3))

toc = [('一', '《论语》今古文版本的源流与分化'),
       ('二', '今文经学与古文经学的学术特征比较'),
       ('三', '《论语》核心概念的今古文诠释差异'),
       ('四', '今古文之争的历史演变脉络'),
       ('五', '今古文之争的当代启示与结论')]
y = Inches(1.6)
for num, t in toc:
    c = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.2), y, Inches(0.6), Inches(0.6))
    c.fill.solid(); c.fill.fore_color.rgb = GOLD; c.line.fill.background()
    p = c.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    _run(p, num, '黑体', 22, True, BG)
    _tb(s, Inches(2.1), y, Inches(9), Inches(0.6), t, size=22, color=WHITE)
    y += Inches(0.95)
notes(s, """这是本次汇报的整体框架。

第一部分，我会梳理《论语》在汉代的版本源流，包括《鲁论》《齐论》和《古论》三个系统，以及它们最终如何融合为通行本。

第二部分，比较今古文两派在文本来源、诠释方法和学术立场上的根本差异。

第三部分，通过"仁""礼"等核心概念的具体诠释案例，展示两派诠释风格的差异。

第四部分，回顾从西汉到清代今古文之争的历史演变。

最后，探讨今古文之争对当代经典诠释和哲学研究的启示。""")

# ================================================================
# SLIDE 2: 秦火与三个版本
# ================================================================
s = blank()
title_bar(s, '一、《论语》今古文版本的源流与分化')
# 左
rect(s, Inches(0.6), Inches(1.5), Inches(5.8), Inches(5.2), BG2)
_tb(s, Inches(0.9), Inches(1.6), Inches(5.2), Inches(0.5),
    '🔥 秦始皇焚书 —— 经典传承的重大危机', name='黑体', size=19, bold=True, color=GOLD)
_mb(s, Inches(0.9), Inches(2.3), Inches(5.2), Inches(4.2), [
    ('秦始皇三十四年（前213年）焚书令，', 15, False, WHITE),
    ('儒家经典遭到系统性破坏', 15, False, WHITE),
    ('', 6),
    ('汉代经典恢复的两条路径：', 17, True, GOLD),
    ('', 4),
    ('① 口耳相传 → 隶书记录 → 今文经学', 16, False, LBLUE),
    ('② 先秦古本 → 古文字书写 → 古文经学', 16, False, LBLUE),
    ('', 6),
    ('《论语》在这两条路径上', 15, False, WHITE),
    ('形成了截然不同的版本系统', 15, False, WHITE),
])
# 右
rect(s, Inches(6.8), Inches(1.5), Inches(5.8), Inches(5.2), BG2)
_tb(s, Inches(7.1), Inches(1.6), Inches(5.2), Inches(0.5),
    '📜 《论语》的三个汉代版本', name='黑体', size=19, bold=True, color=GOLD)
_mb(s, Inches(7.1), Inches(2.3), Inches(5.2), Inches(4.2), [
    ('《鲁论》', 20, True, GOLD),
    ('鲁地儒生传授，20篇，隶书记录', 15, False, WHITE),
    ('汉初流传最广，后为通行本基础', 14, False, LBLUE),
    ('', 8),
    ('《齐论》', 20, True, GOLD),
    ('齐地儒生传授，22篇', 15, False, WHITE),
    ('多出《问王》《知道》两篇', 14, False, LBLUE),
    ('', 8),
    ('《古论》', 20, True, GOLD),
    ('孔子故宅壁中发现，21篇', 15, False, WHITE),
    ('先秦六国古文字书写', 14, False, LBLUE),
])
notes(s, """要理解今古文之争，首先要回到历史现场——秦始皇的焚书令。

公元前213年，秦始皇采纳李斯的建议，下令焚毁民间所藏的《诗》《书》等儒家经典。这一事件对儒家经典的传承造成了毁灭性的打击。

汉代建立之后，经典文献的恢复面临两条路径：

第一条，依靠幸存儒生的记忆。他们通过口耳相传的方式，将经典内容背诵下来，然后用当时通行的隶书记录。这就是"今文"经学的由来。

第二条，依靠偶然发现的先秦古本。最著名的发现就是鲁恭王在孔子故宅的墙壁中发现了用先秦六国古文字书写的《论语》。

《论语》在这两条路径上形成了三个不同的版本：《鲁论》《齐论》和《古论》。

《鲁论》由鲁地儒生传授，共二十篇，是今文系统。《齐论》由齐地儒生传授，共二十二篇，比《鲁论》多出两篇。《古论》则是古文系统，共二十一篇，用先秦古文字书写。""")

# ================================================================
# SLIDE 3: 张侯论
# ================================================================
s = blank()
title_bar(s, '版本的融合与统一 —— "张侯论"')
data = [
    (Inches(0.5), '《鲁论》\n20篇\n（今文·隶书）', RGBColor(0x4A, 0x67, 0x41)),
    (Inches(3.5), '《齐论》\n22篇\n（今文·隶书）', RGBColor(0x4A, 0x67, 0x41)),
    (Inches(6.5), '《古论》\n21篇\n（古文·六国文字）', RGBColor(0x6B, 0x4C, 0x3B)),
]
for x, txt, c in data:
    rect(s, x, Inches(1.8), Inches(2.5), Inches(1.4), c)
    _mb(s, x+Inches(0.1), Inches(1.9), Inches(2.3), Inches(1.2), [(txt, 15, True, WHITE, PP_ALIGN.CENTER)])
_tb(s, Inches(9.2), Inches(2.0), Inches(0.8), Inches(1), '→', name='黑体', size=40, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
rect(s, Inches(10.0), Inches(1.6), Inches(2.8), Inches(2.0), RGBColor(0x8B, 0x69, 0x14))
_mb(s, Inches(10.1), Inches(1.7), Inches(2.6), Inches(1.8), [
    ('张 侯 论', 20, True, WARM, PP_ALIGN.CENTER), ('', 6),
    ('西汉末年 · 张禹编纂', 14, False, WHITE, PP_ALIGN.CENTER), ('', 4),
    ('以《鲁论》为基础', 13, False, LBLUE, PP_ALIGN.CENTER),
    ('参考《齐论》《古论》', 13, False, LBLUE, PP_ALIGN.CENTER), ('', 4),
    ('→ 通行本底本', 14, True, GOLD, PP_ALIGN.CENTER),
])
_mb(s, Inches(1.5), Inches(4.2), Inches(10), Inches(2.5), [
    ('张禹编纂工作的学术意义：', 20, True, GOLD),
    ('• 并非简单的文本拼合，而是比较版本异同后的审慎选择', 16, False, WHITE),
    ('• 体现了今古文经学方法的初步融合', 16, False, WHITE),
    ('• 后来被郑玄用作注释底本，成为《论语》通行本，流传至今', 16, False, WHITE),
    ('• 郑玄兼采今古文，成为东汉经学集大成者', 16, False, LBLUE),
])
notes(s, """这三个版本并不是平行发展的。到了西汉末年，它们走向了融合。

关键人物是张禹。他以《鲁论》为基础，参考《齐论》和《古论》，编纂了一个综合性的《论语》文本，世称"张侯论"。

张禹的工作不是简单地把三个版本拼在一起，而是在仔细比较不同版本的异同之后，做出了审慎的文本选择。这个过程本身，体现了今古文经学方法的某种融合。

"张侯论"后来被郑玄用作注释底本，成为《论语》的通行本，一直流传到今天。郑玄在注释中兼采今古文两派的学术成果，成为东汉经学的集大成者。""")

# ================================================================
# SLIDE 4: 学术特征比较表
# ================================================================
s = blank()
title_bar(s, '二、今文经学 vs 古文经学 —— 学术特征比较')
headers = ['比较维度', '今 文 经 学', '古 文 经 学']
rows = [
    ['文本来源', '汉初儒生口耳相传\n隶书转写', '先秦古本实物发现\n古文字书写'],
    ['核心方法', '微言大义\n经世致用', '训诂考据\n历史还原'],
    ['诠释取向', '政治实用性\n服务现实需要', '学术客观性\n还原经典本义'],
    ['政治地位', '西汉官学，主导朝野', '长期民间学术，东汉兴起'],
    ['代表人物', '董仲舒、公孙弘', '郑玄、许慎、马融'],
]
cw = [Inches(2.0), Inches(4.5), Inches(4.5)]
cs = [Inches(0.8), Inches(2.8), Inches(7.3)]
rh = Inches(0.85); ys = Inches(1.5)
for j, (h, st) in enumerate(zip(headers, cs)):
    rect(s, st, ys, cw[j], Inches(0.5), GOLD)
    _tb(s, st+Inches(0.1), ys+Inches(0.05), cw[j]-Inches(0.2), Inches(0.4),
        h, name='黑体', size=16, bold=True, color=BG, align=PP_ALIGN.CENTER)
for i, row in enumerate(rows):
    y = ys + Inches(0.5) + i*rh
    bg = BG2 if i%2==0 else RGBColor(0x34, 0x45, 0x6B)
    for j, (cell, st) in enumerate(zip(row, cs)):
        rect(s, st, y, cw[j], rh, bg)
        fc = GOLD if j==0 else WHITE
        fb = True if j==0 else False
        fs = 14 if j==0 else 13
        al = PP_ALIGN.LEFT if j==0 else PP_ALIGN.CENTER
        _tb(s, st+Inches(0.1), y+Inches(0.08), cw[j]-Inches(0.2), rh-Inches(0.16),
            cell, size=fs, bold=fb, color=fc, align=al)
notes(s, """现在从学术特征的角度，对今文经学和古文经学做一个系统比较。

这个表格从五个维度展示了两派的差异：

文本来源方面，今文依靠口耳相传后用隶书记录，古文则以先秦古本实物为依据。

核心方法方面，今文强调"微言大义"，古文强调"训诂考据"。

诠释取向上，今文注重政治实用性，古文注重学术客观性。

政治地位上，今文在西汉立为官学，古文长期处于民间，到东汉才兴起。

代表人物方面，今文有董仲舒、公孙弘，古文有郑玄、许慎、马融。""")

# ================================================================
# SLIDE 5: 诠释方法深度比较
# ================================================================
s = blank()
title_bar(s, '诠释方法的核心分歧')
# 今文
rect(s, Inches(0.5), Inches(1.5), Inches(5.8), Inches(5.5), BG2)
_tb(s, Inches(0.8), Inches(1.6), Inches(5.2), Inches(0.5),
    '今文经学："微言大义"', name='黑体', size=22, bold=True, color=GOLD)
_mb(s, Inches(0.8), Inches(2.3), Inches(5.2), Inches(4.5), [
    ('方法论核心：', 17, True, GOLD),
    ('经典中蕴含圣人的深层政治意图', 15, False, WHITE),
    ('和道德理想，诠释者需揭示"大义"', 15, False, WHITE),
    ('', 6), ('特点：', 17, True, GOLD),
    ('• 强调经典的现实关怀', 14, False, WHITE),
    ('• 诠释服务于政治需要', 14, False, WHITE),
    ('• 赋予经典鲜活的生命力', 14, False, WHITE),
    ('', 6), ('典型体现：', 17, True, GOLD),
    ('董仲舒"天人感应"说', 15, False, LBLUE),
    ('将《春秋》诠释为维护大一统', 14, False, LBLUE),
    ('的理论基础', 14, False, LBLUE),
    ('', 6), ('⚠ 风险：诠释主观化、政治化', 14, True, RED),
])
# 古文
rect(s, Inches(6.8), Inches(1.5), Inches(5.8), Inches(5.5), BG2)
_tb(s, Inches(7.1), Inches(1.6), Inches(5.2), Inches(0.5),
    '古文经学："训诂考据"', name='黑体', size=22, bold=True, color=GOLD)
_mb(s, Inches(7.1), Inches(2.3), Inches(5.2), Inches(4.5), [
    ('方法论核心：', 17, True, GOLD),
    ('通过文字、音韵、训诂的精细研究，', 15, False, WHITE),
    ('还原经典文本的原始含义', 15, False, WHITE),
    ('', 6), ('特点：', 17, True, GOLD),
    ('• 强调历史还原和学术客观性', 14, False, WHITE),
    ('• 诠释者首要任务是理解本义', 14, False, WHITE),
    ('• 避免主观意志强加于经典', 14, False, WHITE),
    ('', 6), ('典型体现：', 17, True, GOLD),
    ('郑玄《论语注》', 15, False, LBLUE),
    ('许慎《说文解字》', 15, False, LBLUE),
    ('马融古文经传', 15, False, LBLUE),
    ('', 6), ('⚠ 风险：陷入繁琐考证、脱离现实', 14, True, RED),
])
notes(s, """诠释方法的差异，是今古文之争最核心的分歧。

左边是今文经学的"微言大义"。这种方法认为，经典文本中蕴含着圣人的深层政治意图和道德理想。诠释者不是停留在字面意思上，而是要通过解读"微言"来揭示"大义"。

董仲舒的"天人感应"说就是最典型体现。他把《春秋》中的灾异诠释为上天的警示，为汉武帝的大一统政治提供理论基础。

右边是古文经学的"训诂考据"。这种方法强调通过文字、音韵、训诂的精细研究来还原本义。郑玄的《论语注》、许慎的《说文解字》都是杰出代表。

两种方法各有优劣。今文赋予经典强烈现实意义但可能过度解读。古文保证客观性但可能使研究过于琐碎。""")

# ================================================================
# SLIDE 6: 个案分析
# ================================================================
s = blank()
title_bar(s, '三、核心概念的今古文诠释差异 —— 个案分析')
rect(s, Inches(0.5), Inches(1.5), Inches(5.8), Inches(2.3), BG2)
_tb(s, Inches(0.8), Inches(1.6), Inches(5.2), Inches(0.5),
    '"仁"的诠释差异', name='黑体', size=20, bold=True, color=GOLD)
_mb(s, Inches(0.8), Inches(2.2), Inches(5.2), Inches(1.5), [
    ('今文经学：', 16, True, LBLUE),
    ('将"仁"与政治治理和社会秩序相联系', 14, False, WHITE),
    ('强调"仁政"的政治功能', 14, False, WHITE),
    ('', 5), ('古文经学：', 16, True, LBLUE),
    ('从文字训诂角度还原"仁"的原始语义', 14, False, WHITE),
    ('强调其作为道德品质的内涵', 14, False, WHITE),
])
rect(s, Inches(6.8), Inches(1.5), Inches(5.8), Inches(2.3), BG2)
_tb(s, Inches(7.1), Inches(1.6), Inches(5.2), Inches(0.5),
    '"礼"的诠释差异', name='黑体', size=20, bold=True, color=GOLD)
_mb(s, Inches(7.1), Inches(2.2), Inches(5.2), Inches(1.5), [
    ('今文经学：', 16, True, LBLUE),
    ('将"礼"理解为维系社会等级秩序的', 14, False, WHITE),
    ('制度规范，强调政治功能', 14, False, WHITE),
    ('', 5), ('古文经学：', 16, True, LBLUE),
    ('将"礼"视为文化传统和行为规范', 14, False, WHITE),
    ('注重历史演变和文化意义', 14, False, WHITE),
])
rect(s, Inches(0.5), Inches(4.2), Inches(12.1), Inches(2.6), BG2)
_tb(s, Inches(0.8), Inches(4.3), Inches(11), Inches(0.5),
    '诠释风格对比', name='黑体', size=19, bold=True, color=GOLD)
_mb(s, Inches(0.8), Inches(4.9), Inches(5.5), Inches(1.8), [
    ('今文经学诠释风格：', 16, True, LBLUE),
    ('• 强烈的目的论色彩', 14, False, WHITE),
    ('• 带着政治或道德目的解读经典', 14, False, WHITE),
    ('✅ 优势：赋予经典强烈现实意义', 13, False, GREEN),
    ('❌ 缺陷：可能导致过度解读或曲解', 13, False, RED),
])
_mb(s, Inches(6.8), Inches(4.9), Inches(5.5), Inches(1.8), [
    ('古文经学诠释风格：', 16, True, LBLUE),
    ('• 审慎和客观的学术态度', 14, False, WHITE),
    ('• 力求还原经典文本原始含义', 14, False, WHITE),
    ('✅ 优势：保证学术客观性和可靠性', 13, False, GREEN),
    ('❌ 缺陷：可能过于琐碎、脱离现实', 13, False, RED),
])
notes(s, """来看两个具体个案——《论语》中两个最核心的概念："仁"和"礼"。

先看"仁"。今文经学家倾向于将"仁"与政治治理和社会秩序相联系，强调"仁政"的政治功能。古文经学家则从文字训诂角度，还原"仁"作为道德品质的内涵。

再看"礼"。今文经学家将"礼"理解为维系社会等级秩序的制度规范。古文经学家则将"礼"视为文化传统和行为规范，注重其历史演变和文化意义。

两种诠释风格各有优劣。今文具有目的论色彩，赋予经典现实意义但可能过度解读。古文审慎客观，保证可靠性但可能过于琐碎。""")

# ================================================================
# SLIDE 7: 历史演变时间线
# ================================================================
s = blank()
title_bar(s, '四、今古文之争的历史演变脉络')
periods = [
    ('西汉', '今文经学鼎盛',
     ['汉武帝"罢黜百家，独尊儒术"',
      '今文经学立为官学，主导朝野',
      '《鲁论》《齐论》广泛流传',
      '董仲舒"天人感应"影响深远'],
     RGBColor(0x4A, 0x67, 0x41)),
    ('东汉', '古文经学兴起',
     ['古文经典不断发现',
      '郑玄兼采今古文，集大成',
      '以"张侯论"为底本作注',
      '古文经学影响力日益增强'],
     RGBColor(0x6B, 0x4C, 0x3B)),
    ('清代', '今古文之争复兴',
     ['乾嘉考据学大兴',
      '庄存与、刘逢禄复兴今文',
      '康有为《新学伪经考》',
      '今文经学与变法维新结合'],
     RGBColor(0x4A, 0x4A, 0x7A)),
]
xpos = [Inches(0.5), Inches(4.6), Inches(8.7)]
for idx, (period, title, items, color) in enumerate(periods):
    x = xpos[idx]
    rect(s, x, Inches(1.5), Inches(3.6), Inches(0.65), color)
    _tb(s, x, Inches(1.55), Inches(3.6), Inches(0.55),
        f'{period}：{title}', name='黑体', size=17, bold=True, color=WARM, align=PP_ALIGN.CENTER)
    if idx < 2:
        gold_bar(s, x+Inches(3.6), Inches(1.8), xpos[idx+1]-x-Inches(3.6))
    y = Inches(2.4)
    for item in items:
        _tb(s, x+Inches(0.2), y, Inches(3.2), Inches(0.6), '• '+item, size=13, color=WHITE)
        y += Inches(0.85)
rect(s, Inches(0.5), Inches(6.0), Inches(12.3), Inches(0.7), RGBColor(0x3D, 0x2B, 0x1F))
_tb(s, Inches(1.0), Inches(6.05), Inches(11.3), Inches(0.6),
    '每一次学术转向都深刻影响了中国思想文化的发展走向', name='黑体', size=17, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
notes(s, """今古文之争贯穿了整个中国经学史，分为三个关键阶段。

西汉时期，今文经学独尊。汉武帝"罢黜百家，独尊儒术"，所尊之儒学就是今文经学。

东汉时期，古文经学兴起。郑玄是这个阶段的关键人物，他兼采今古文，成为东汉经学的集大成者。

清代，今古文之争再次复兴。乾嘉时期考据学大兴，但到晚清，庄存与、刘逢禄重新强调今文经学，康有为更是将今文经学与变法维新相结合。

每一次学术转向都深刻影响了中国思想文化的发展走向。""")

# ================================================================
# SLIDE 8: 康有为
# ================================================================
s = blank()
title_bar(s, '特别关注：康有为与晚清今古文之争')
_mb(s, Inches(0.8), Inches(1.6), Inches(11.5), Inches(5.0), [
    ('康有为的"伪经"说', 24, True, GOLD), ('', 8),
    ('核心观点：古文经乃刘歆为配合王莽篡汉而伪造', 18, False, WHITE),
    ('', 8), ('代表著作：', 19, True, LBLUE),
    ('• 《新学伪经考》—— 系统论证古文经为伪经', 16, False, WHITE),
    ('• 《孔子改制考》—— 将孔子塑造为改革家', 16, False, WHITE),
    ('', 8), ('政治动机：为变法维新提供理论支持', 17, False, WHITE),
    ('', 8), ('学术影响：', 19, True, LBLUE),
    ('• 学术上存在诸多争议，至今未有定论', 16, False, WHITE),
    ('• 但在当时社会政治环境下产生了巨大影响', 16, False, WHITE),
    ('• 今文经学从纯学术辩论转化为政治变革的思想武器', 16, False, LBLUE),
])
notes(s, """在晚清今古文之争的复兴中，康有为是绕不开的人物。

康有为提出了"伪经"说——认为古文经乃刘歆为配合王莽篡汉而伪造。这直接动摇了古文经学的根基。

他写了两本影响深远的著作：《新学伪经考》系统论证古文经是伪经；《孔子改制考》把孔子塑造成改革家，为变法维新提供历史合法性。

这些观点在学术上有很多争议。但在当时的社会政治环境下，这些观点产生了巨大影响。

康有为成功地将今文经学从纯学术辩论转化为政治变革的思想武器。这也印证了——今文经学的"微言大义"方法本身就具有强烈的政治实用性。""")

# ================================================================
# SLIDE 9: 当代启示
# ================================================================
s = blank()
title_bar(s, '五、今古文之争的当代启示')
_mb(s, Inches(0.8), Inches(1.6), Inches(5.5), Inches(5.0), [
    ('两种诠释范式', 22, True, GOLD), ('', 6),
    ('以意义阐释为核心', 17, False, LBLUE),
    ('→ 强调经典的现实关怀和时代价值', 15, False, WHITE),
    ('', 8), ('以文本还原为核心', 17, False, LBLUE),
    ('→ 强调经典的原始含义和历史语境', 15, False, WHITE),
    ('', 10), ('当代经典研究应当超越今古文', 18, True, GOLD),
    ('之争的简单对立，将两种范式', 18, True, GOLD),
    ('有机结合起来', 18, True, GOLD),
])
_mb(s, Inches(6.8), Inches(1.6), Inches(5.8), Inches(5.0), [
    ('对当代中国哲学研究的意义', 22, True, GOLD), ('', 6),
    ('① 关注经典诠释的方法论问题', 16, False, WHITE),
    ('      避免将单一诠释方法绝对化', 14, False, LBLUE),
    ('', 5), ('② 反思经典诠释与政治社会的复杂关系', 16, False, WHITE),
    ('      思考学术研究的社会功能', 14, False, LBLUE),
    ('', 5), ('③ 发掘中国传统学术中的诠释资源', 16, False, WHITE),
    ('      为当代中国哲学创新提供文化根基', 14, False, LBLUE),
    ('', 5), ('④ 推动中国传统经典的当代阐释', 16, False, WHITE),
    ('      走向新的深度与高度', 14, False, LBLUE),
])
notes(s, """今古文之争对当代仍有重要启示。

从当代学术视角看，今古文之争反映了经典诠释中的两种基本范式：一种以意义阐释为核心，强调现实关怀和时代价值；另一种以文本还原为核心，强调原始含义和历史语境。

这两种范式各有合理性和局限性。当代的经典研究应当超越简单对立，将两种范式有机结合。

具体来说有四个方面的重要意义：第一，关注经典诠释的方法论问题；第二，反思经典诠释与政治社会的复杂关系；第三，发掘中国传统学术中丰富的诠释资源；第四，推动中国传统经典的当代阐释走向新的深度与高度。""")

# ================================================================
# SLIDE 10: 结论
# ================================================================
s = blank()
fill_bg(s, BG)
rect(s, Inches(0), Inches(2.8), Inches(13.333), Inches(0.06), GOLD)
rect(s, Inches(0), Inches(5.8), Inches(13.333), Inches(0.04), GOLD2)
_tb(s, Inches(1.5), Inches(0.8), Inches(10.3), Inches(1.0),
    '结  论', name='黑体', size=34, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
_mb(s, Inches(1.5), Inches(3.3), Inches(10.3), Inches(2.2), [
    ('今古文之争并非简单的文本真伪之辩，', 22, True, WHITE, PP_ALIGN.CENTER),
    ('而是两种经典诠释范式的根本对立', 22, True, WHITE, PP_ALIGN.CENTER),
    ('', 8),
    ('从西汉的今文独尊 → 东汉的古文兴起 → 清代的今古文复兴', 18, False, LBLUE, PP_ALIGN.CENTER),
    ('', 8),
    ('在当代学术语境中，今古文之争的理论价值更加凸显。', 19, False, WHITE, PP_ALIGN.CENTER),
    ('超越简单对立，创造性融合两种诠释范式，', 19, False, WHITE, PP_ALIGN.CENTER),
    ('推动中国传统经典的当代阐释走向新的深度与高度。', 19, False, WHITE, PP_ALIGN.CENTER),
])
_tb(s, Inches(1.5), Inches(6.1), Inches(10.3), Inches(0.6),
    '感谢聆听，敬请批评指正！', name='黑体', size=22, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
notes(s, f"""最后，我来做一个总结。

《论语》今文经学与古文经学之争，是中国经学史上最为重要的学术争论之一。

从版本源流来看，今文以口耳相传、隶书记录为特征，古文以先秦古本、古文字书写为依据。

从学术特征来看，今文注重微言大义与经世致用，古文强调训诂考据与历史还原。

从诠释方法来看，今文具有强烈的政治目的论色彩，古文更加注重学术的客观性和独立性。

今古文之争并非简单的文本真伪之辩，而是两种经典诠释范式的根本对立。

以上就是我的汇报。感谢各位的聆听，敬请批评指正。

——{AUTHOR}，《{COURSE}》课程论文汇报""")

# ================================================================
# SAVE
# ================================================================
ppt_path = '/home/admin/.openclaw/workspace/papers/lunyu-jingu/ppt/论语今古文经学比较研究.pptx'
prs.save(ppt_path)
print(f'✅ PPT: {ppt_path}')
print(f'   幻灯片: {len(prs.slides)}')

# 讲稿
md = [f'# 《论语今文经学与古文经学比较研究》演讲讲稿\n',
      f'**作者：{AUTHOR}**  |  **课程：{COURSE}**  |  **导师：{ADVISOR}**  |  **日期：{date.today().strftime("%Y年%-m月%-d日")}**\n',
      '---\n']
for i, sl in enumerate(prs.slides):
    n = sl.notes_slide.notes_text_frame.text if sl.has_notes_slide else '（无讲稿）'
    if i == 0: title = '封面'
    elif i == 1: title = '汇报提纲'
    elif i == len(prs.slides)-1: title = '结论'
    else:
        title = ''
        for shape in sl.shapes:
            if shape.has_text_frame:
                t = shape.text_frame.text.strip()
                if t and len(t) < 50: title = t; break
    md.append(f'\n## 第{i+1}页：{title}\n\n{n}\n\n---\n')
gp = '/home/admin/.openclaw/workspace/papers/lunyu-jingu/ppt/讲稿.md'
with open(gp, 'w', encoding='utf-8') as f:
    f.write('\n'.join(md))
print(f'✅ 讲稿: {gp}')
