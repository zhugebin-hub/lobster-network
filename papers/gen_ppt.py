#!/usr/bin/env python3
"""AI工具使用指南 PPT - 科技风 22页"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height

# 配色
BG = RGBColor(0x0A, 0x0E, 0x2A)
CARD = RGBColor(0x16, 0x1D, 0x50)
CYAN = RGBColor(0x00, 0xD4, 0xFF)
MAGENTA = RGBColor(0xE0, 0x40, 0xFB)
GREEN = RGBColor(0x00, 0xE6, 0x76)
YELLOW = RGBColor(0xFF, 0xEA, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0xCC, 0xCC, 0xDD)
ORANGE = RGBColor(0xFF, 0x91, 0x00)
RED = RGBColor(0xFF, 0x52, 0x52)
TEAL = RGBColor(0x64, 0xFF, 0xDA)
DIM = RGBColor(0x00, 0x99, 0xBB)
DARK = RGBColor(0x1A, 0x22, 0x55)

def bg(slide, c=BG):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = c

def rect(slide, l, t, w, h, fc, bc=None, bw=Pt(0)):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fc
    if bc: s.line.fill.solid(); s.line.fill.fore_color.rgb = bc; s.line.width = bw
    else: s.line.fill.background()
    return s

def rrect(slide, l, t, w, h, fc, bc=None, bw=Pt(1)):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fc
    if bc: s.line.fill.solid(); s.line.fill.fore_color.rgb = bc; s.line.width = bw
    else: s.line.fill.background()
    return s

def circ(slide, l, t, sz, fc):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, l, t, sz, sz)
    s.fill.solid(); s.fill.fore_color.rgb = fc
    s.line.fill.background()
    return s

def hline(slide, l, t, w, c=CYAN, bw=Pt(2)):
    return rect(slide, l, t, w, bw, c)

def vline(slide, l, t, h, c=CYAN, bw=Pt(3)):
    return rect(slide, l, t, bw, h, c)

def tb(slide, l, t, w, h, text, sz=18, c=WHITE, b=False, a=PP_ALIGN.LEFT, fn='微软雅黑'):
    tx = slide.shapes.add_textbox(l, t, w, h)
    tf = tx.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(sz); p.font.color.rgb = c; p.font.bold = b; p.font.name = fn
    p.alignment = a
    return tx

def mtb(slide, l, t, w, h, lines, dsz=16, ls=1.3):
    tx = slide.shapes.add_textbox(l, t, w, h)
    tf = tx.text_frame; tf.word_wrap = True
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if isinstance(ln, str):
            p.text = ln; p.font.size = Pt(dsz); p.font.color.rgb = WHITE
        else:
            p.text = ln.get('text','')
            p.font.size = Pt(ln.get('size',dsz))
            p.font.color.rgb = ln.get('color',WHITE)
            p.font.bold = ln.get('bold',False)
            if 'align' in ln: p.alignment = ln['align']
        p.line_spacing = Pt(int(p.font.size.pt * ls)) if p.font.size else Pt(int(dsz*ls))
    return tx

def decor(slide):
    hline(slide, Inches(0), Inches(0), Inches(3), CYAN, Pt(3))
    vline(slide, Inches(0), Inches(0), Inches(2), CYAN, Pt(3))
    hline(slide, Inches(10.333), Inches(7.2), Inches(3), CYAN, Pt(3))
    vline(slide, Inches(13.033), Inches(5.5), Inches(2), CYAN, Pt(3))
    circ(slide, Inches(11.5), Inches(0.3), Inches(0.15), CYAN)
    circ(slide, Inches(11.8), Inches(0.5), Inches(0.1), MAGENTA)
    circ(slide, Inches(0.3), Inches(6.8), Inches(0.12), CYAN)

def hdr(slide, n, title):
    tb(slide, Inches(0.8), Inches(0.3), Inches(3), Inches(0.5),
       f"0{n}", sz=24, c=CYAN, b=True)
    tb(slide, Inches(3.5), Inches(0.3), Inches(8), Inches(0.5),
       title, sz=20, c=GRAY)
    hline(slide, Inches(0.8), Inches(0.85), Inches(11.7), CYAN, Pt(1))

def pg(slide, n, total=22):
    tb(slide, Inches(12.2), Inches(7.0), Inches(1), Inches(0.4),
       f"{n}/{total}", sz=12, c=DIM, a=PP_ALIGN.RIGHT)

# ========== SLIDE 1: COVER ==========
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
circ(s, Inches(-2), Inches(-2), Inches(6), RGBColor(0x0D,0x14,0x3A))
circ(s, Inches(10), Inches(4), Inches(5), RGBColor(0x0D,0x14,0x3A))
circ(s, Inches(5), Inches(-1), Inches(3), RGBColor(0x0F,0x18,0x45))
hline(s, Inches(1.5), Inches(2.8), Inches(10.3), CYAN, Pt(2))
hline(s, Inches(1.5), Inches(5.2), Inches(10.3), MAGENTA, Pt(1))
tb(s, Inches(1.5), Inches(3.0), Inches(10.3), Inches(1.2),
   "AI工具使用指南", sz=54, c=WHITE, b=True, a=PP_ALIGN.CENTER)
tb(s, Inches(1.5), Inches(4.3), Inches(10.3), Inches(0.6),
   "中职高一信息技术课 · 认识身边的AI助手", sz=24, c=CYAN, a=PP_ALIGN.CENTER)
tb(s, Inches(1.5), Inches(5.5), Inches(10.3), Inches(0.5),
   "5大AI工具详解 ｜ 实战案例 ｜ 使用步骤 ｜ 优缺点分析",
   sz=16, c=GRAY, a=PP_ALIGN.CENTER)
circ(s, Inches(1.8), Inches(5.6), Inches(0.1), CYAN)
circ(s, Inches(11.4), Inches(5.6), Inches(0.1), MAGENTA)
pg(s, 1)

# ========== SLIDE 2: TOC ==========
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); decor(s)
tb(s, Inches(0.8), Inches(0.3), Inches(5), Inches(0.6), "目  录", sz=36, c=WHITE, b=True)
hline(s, Inches(0.8), Inches(0.9), Inches(4), CYAN, Pt(2))
y = Inches(1.3)
for num, title, color in [
    ("01","什么是AI工具？为什么学？",CYAN),
    ("02","AI工具全景概览（5大工具）",CYAN),
    ("03","豆包 — 全能AI对话助手",MAGENTA),
    ("04","DeepSeek — 深度思考AI",TEAL),
    ("05","醒图 — AI智能修图神器",GREEN),
    ("06","通义千问 — 阿里AI助手",ORANGE),
    ("07","剪映 — AI视频创作工具",RED),
    ("08","工具对比与选择建议",YELLOW),
    ("09","AI工具使用注意事项",RED),
    ("10","课堂实践任务",GREEN),
]:
    tb(s, Inches(1.0), y, Inches(0.8), Inches(0.4), num, sz=20, c=color, b=True)
    tb(s, Inches(1.8), y, Inches(8), Inches(0.4), title, sz=18, c=WHITE)
    hline(s, Inches(1.0), y+Inches(0.42), Inches(10.5), DARK, Pt(1))
    y += Inches(0.55)
pg(s, 2)

# ========== SLIDE 3: WHAT IS AI ==========
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); decor(s); hdr(s,1,"什么是AI工具？")
rrect(s, Inches(0.8), Inches(1.2), Inches(5.5), Inches(5.5), CARD, CYAN, Pt(1))
tb(s, Inches(1.0), Inches(1.4), Inches(5), Inches(0.5), "🤖 什么是AI工具？", sz=22, c=CYAN, b=True)
hline(s, Inches(1.0), Inches(1.9), Inches(5), DIM, Pt(1))
mtb(s, Inches(1.0), Inches(2.1), Inches(5), Inches(4), [
    "AI工具 = 人工智能技术驱动的",
    "软件/平台/应用",
    "",
    "它们能模拟人类的：",
    "  ✓ 语言理解与生成",
    "  ✓ 图像识别与创作",
    "  ✓ 逻辑推理与分析",
    "  ✓ 语音识别与合成",
    "",
    "简单说：AI工具就是",
    "你的「数字助手」！",
], dsz=16)

rrect(s, Inches(6.8), Inches(1.2), Inches(5.5), Inches(5.5), CARD, MAGENTA, Pt(1))
tb(s, Inches(7.0), Inches(1.4), Inches(5), Inches(0.5), "🎯 为什么中职生要学AI？", sz=22, c=MAGENTA, b=True)
hline(s, Inches(7.0), Inches(1.9), Inches(5), RGBColor(0xAA,0x30,0xCC), Pt(1))
mtb(s, Inches(7.0), Inches(2.1), Inches(5), Inches(4), [
    {"text":"1. 未来职场必备技能","size":15,"color":WHITE,"bold":True},
    {"text":"   → 90%岗位将与AI协作","size":14,"color":GRAY},
    {"text":"", "size":6},
    {"text":"2. 提升学习效率","size":15,"color":WHITE,"bold":True},
    {"text":"   → 快速查找资料、解答疑问","size":14,"color":GRAY},
    {"text":"", "size":6},
    {"text":"3. 激发创意思维","size":15,"color":WHITE,"bold":True},
    {"text":"   → AI帮你突破灵感瓶颈","size":14,"color":GRAY},
    {"text":"", "size":6},
    {"text":"4. 增强就业竞争力","size":15,"color":WHITE,"bold":True},
    {"text":"   → 掌握AI = 掌握未来","size":14,"color":GRAY},
], ls=1.3)
pg(s, 3)

# ========== SLIDE 4: OVERVIEW ==========
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); decor(s); hdr(s,2,"AI工具全景概览")
tools = [
    ("📱 豆包","字节跳动","AI对话/写作/问答",CYAN),
    ("🧠 DeepSeek","深度求索","深度推理/代码/分析",TEAL),
    ("🎨 醒图","南京代码","AI修图/设计/美颜",GREEN),
    ("💬 通义千问","阿里巴巴","对话/创作/多模态",ORANGE),
    ("🎬 剪映","字节跳动","AI视频剪辑/特效",RED),
]
x = Inches(0.6)
for name, co, desc, color in tools:
    rrect(s, x, Inches(1.3), Inches(2.3), Inches(5.2), CARD, color, Pt(2))
    circ(s, x+Inches(0.7), Inches(1.6), Inches(0.9), color)
    tb(s, x+Inches(0.15), Inches(2.6), Inches(2.0), Inches(0.8),
       name, sz=22, c=color, b=True, a=PP_ALIGN.CENTER)
    tb(s, x+Inches(0.15), Inches(3.5), Inches(2.0), Inches(0.4),
       f"开发商：{co}", sz=12, c=GRAY, a=PP_ALIGN.CENTER)
    hline(s, x+Inches(0.3), Inches(4.0), Inches(1.7), color, Pt(1))
    tb(s, x+Inches(0.15), Inches(4.2), Inches(2.0), Inches(0.4),
       "核心功能：", sz=12, c=DIM, a=PP_ALIGN.CENTER)
    tb(s, x+Inches(0.15), Inches(4.5), Inches(2.0), Inches(0.8),
       desc, sz=14, c=WHITE, a=PP_ALIGN.CENTER)
    tb(s, x+Inches(0.15), Inches(5.5), Inches(2.0), Inches(0.4),
       "免费使用 ✓", sz=12, c=GREEN, a=PP_ALIGN.CENTER)
    x += Inches(2.5)
pg(s, 4)

# ========== SLIDE 5-6: DOUBAO ==========
# Slide 5
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); decor(s); hdr(s,3,"豆包 — 字节跳动AI助手")
rrect(s, Inches(0.8), Inches(1.2), Inches(5.8), Inches(5.8), CARD, CYAN, Pt(2))
tb(s, Inches(1.0), Inches(1.4), Inches(5), Inches(0.5), "📱 豆包", sz=28, c=CYAN, b=True)
hline(s, Inches(1.0), Inches(1.9), Inches(5), CYAN, Pt(1))
mtb(s, Inches(1.0), Inches(2.1), Inches(5.3), Inches(4.5), [
    {"text":"开发商：字节跳动（抖音母公司）","size":15,"color":WHITE,"bold":True},
    {"text":"","size":8},
    {"text":"平台：网页版 / APP（iOS+Android）","size":14,"color":GRAY},
    {"text":"网址：www.doubao.com","size":14,"color":CYAN},
    {"text":"","size":8},
    {"text":"核心功能：","size":15,"color":CYAN,"bold":True},
    {"text":"  • 智能对话问答","size":13,"color":GRAY},
    {"text":"  • 文章写作/改写","size":13,"color":GRAY},
    {"text":"  • 知识讲解/答疑","size":13,"color":GRAY},
    {"text":"  • 图片生成（文生图）","size":13,"color":GRAY},
    {"text":"  • 语言翻译","size":13,"color":GRAY},
    {"text":"  • 文件分析（上传文档）","size":13,"color":GRAY},
], ls=1.3)

rrect(s, Inches(7.0), Inches(1.2), Inches(5.5), Inches(5.8), CARD, GREEN, Pt(2))
tb(s, Inches(7.2), Inches(1.4), Inches(5), Inches(0.5), "💡 实用案例", sz=22, c=GREEN, b=True)
hline(s, Inches(7.2), Inches(1.9), Inches(5), GREEN, Pt(1))
mtb(s, Inches(7.2), Inches(2.1), Inches(5), Inches(4.5), [
    {"text":"案例1：写作业助手","size":15,"color":YELLOW,"bold":True},
    {"text":"学生：\"帮我解释光合作用，","size":12,"color":GRAY},
    {"text":"  用通俗语言\"","size":12,"color":GRAY},
    {"text":"豆包：生成清晰易懂的","size":12,"color":GRAY},
    {"text":"  分步骤解释 ✓","size":12,"color":GREEN},
    {"text":"","size":6},
    {"text":"案例2：作文润色","size":15,"color":YELLOW,"bold":True},
    {"text":"学生：把这段文字改写","size":12,"color":GRAY},
    {"text":"得更生动，加入修辞手法","size":12,"color":GRAY},
    {"text":"豆包：自动优化表达，","size":12,"color":GRAY},
    {"text":"  提升文采 ✓","size":12,"color":GREEN},
    {"text":"","size":6},
    {"text":"案例3：职业规划","size":15,"color":YELLOW,"bold":True},
    {"text":"学生：\"中职电商专业","size":12,"color":GRAY},
    {"text":"  毕业后可以做什么？\"","size":12,"color":GRAY},
    {"text":"豆包：列出就业方向+技能","size":12,"color":GRAY},
    {"text":"  建议+薪资参考 ✓","size":12,"color":GREEN},
], ls=1.2)
pg(s, 5)

# Slide 6
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); decor(s); hdr(s,3,"豆包 — 优缺点与使用步骤")
rrect(s, Inches(0.8), Inches(1.2), Inches(3.5), Inches(5.8), CARD, GREEN, Pt(2))
tb(s, Inches(1.0), Inches(1.4), Inches(3), Inches(0.5), "✅ 优点", sz=22, c=GREEN, b=True)
hline(s, Inches(1.0), Inches(1.9), Inches(3), GREEN, Pt(1))
mtb(s, Inches(1.0), Inches(2.1), Inches(3), Inches(4.5),
    [{"text":p,"size":14,"color":WHITE} for p in [
        "• 免费使用，无门槛","• 中文理解能力强","• 响应速度快",
        "• 支持多种格式输入","• 与抖音生态联动","• 适合日常学习问答",
        "• 界面简洁易用","• 支持图片生成",
    ]], ls=1.5)

rrect(s, Inches(4.7), Inches(1.2), Inches(3.5), Inches(5.8), CARD, RED, Pt(2))
tb(s, Inches(4.9), Inches(1.4), Inches(3), Inches(0.5), "⚠️ 缺点", sz=22, c=RED, b=True)
hline(s, Inches(4.9), Inches(1.9), Inches(3), RED, Pt(1))
mtb(s, Inches(4.9), Inches(2.1), Inches(3), Inches(4.5),
    [{"text":c,"size":14,"color":WHITE} for c in [
        "• 深度推理能力有限","• 长文本可能遗漏",
        "• 专业领域知识不足","• 可能产生\"幻觉\"",
        "  （编造信息）","• 不能替代独立思考",
        "• 免费额度有限制","• 依赖网络连接",
    ]], ls=1.5)

rrect(s, Inches(8.6), Inches(1.2), Inches(3.9), Inches(5.8), CARD, CYAN, Pt(2))
tb(s, Inches(8.8), Inches(1.4), Inches(3.5), Inches(0.5), "📋 使用步骤", sz=22, c=CYAN, b=True)
hline(s, Inches(8.8), Inches(1.9), Inches(3.5), CYAN, Pt(1))
mtb(s, Inches(8.8), Inches(2.1), Inches(3.5), Inches(4.5), [
    {"text":"Step 1：下载/访问","size":15,"color":YELLOW,"bold":True},
    {"text":"  手机应用商店搜\"豆包\"","size":12,"color":GRAY},
    {"text":"  或网页访问 doubao.com","size":12,"color":GRAY},
    {"text":"","size":6},
    {"text":"Step 2：注册登录","size":15,"color":YELLOW,"bold":True},
    {"text":"  手机号/抖音账号登录","size":12,"color":GRAY},
    {"text":"","size":6},
    {"text":"Step 3：输入问题","size":15,"color":YELLOW,"bold":True},
    {"text":"  在对话框输入你的问题","size":12,"color":GRAY},
    {"text":"  问题越具体，回答越好","size":12,"color":GRAY},
    {"text":"","size":6},
    {"text":"Step 4：追问优化","size":15,"color":YELLOW,"bold":True},
    {"text":"  对回答不满意可继续追问","size":12,"color":GRAY},
    {"text":"","size":6},
    {"text":"Step 5：复制使用","size":15,"color":YELLOW,"bold":True},
    {"text":"  复制有用内容到自己的","size":12,"color":GRAY},
    {"text":"  文档中整理使用","size":12,"color":GRAY},
], ls=1.2)
pg(s, 6)

# ========== SLIDE 7-8: DEEPSEEK ==========
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); decor(s); hdr(s,4,"DeepSeek — 深度推理AI")
rrect(s, Inches(0.8), Inches(1.2), Inches(5.8), Inches(5.8), CARD, TEAL, Pt(2))
tb(s, Inches(1.0), Inches(1.4), Inches(5), Inches(0.5), "🧠 DeepSeek", sz=28, c=TEAL, b=True)
hline(s, Inches(1.0), Inches(1.9), Inches(5), TEAL, Pt(1))
mtb(s, Inches(1.0), Inches(2.1), Inches(5.3), Inches(4.5), [
    {"text":"开发商：深度求索（中国AI公司）","size":15,"color":WHITE,"bold":True},
    {"text":"","size":8},
    {"text":"平台：网页版 / APP","size":14,"color":GRAY},
    {"text":"网址：www.deepseek.com","size":14,"color":TEAL},
    {"text":"","size":8},
    {"text":"核心特色：","size":15,"color":TEAL,"bold":True},
    {"text":"  • 深度推理能力强大","size":13,"color":GRAY},
    {"text":"  • 数学/逻辑题擅长","size":13,"color":GRAY},
    {"text":"  • 代码编写与调试","size":13,"color":GRAY},
    {"text":"  • 长文本分析处理","size":13,"color":GRAY},
    {"text":"  • 多语言翻译精准","size":13,"color":GRAY},
    {"text":"  • 免费额度充足","size":13,"color":GRAY},
], ls=1.3)

rrect(s, Inches(7.0), Inches(1.2), Inches(5.5), Inches(5.8), CARD, CYAN, Pt(2))
tb(s, Inches(7.2), Inches(1.4), Inches(5), Inches(0.5), "💡 实用案例", sz=22, c=CYAN, b=True)
hline(s, Inches(7.2), Inches(1.9), Inches(5), CYAN, Pt(1))
mtb(s, Inches(7.2), Inches(2.1), Inches(5), Inches(4.5), [
    {"text":"案例1：数学题详解","size":15,"color":YELLOW,"bold":True},
    {"text":"学生：\"解方程 2x²+5x-3=0","size":12,"color":GRAY},
    {"text":"  请给出详细步骤\"","size":12,"color":GRAY},
    {"text":"DeepSeek：逐步推导，","size":12,"color":GRAY},
    {"text":"  解释每一步原理 ✓","size":12,"color":GREEN},
    {"text":"","size":6},
    {"text":"案例2：编程入门","size":15,"color":YELLOW,"bold":True},
    {"text":"学生：\"用Python写一个","size":12,"color":GRAY},
    {"text":"  猜数字小游戏\"","size":12,"color":GRAY},
    {"text":"DeepSeek：生成完整代码+","size":12,"color":GRAY},
    {"text":"  逐行注释讲解 ✓","size":12,"color":GREEN},
    {"text":"","size":6},
    {"text":"案例3：长文档分析","size":15,"color":YELLOW,"bold":True},
    {"text":"学生：上传一篇3000字","size":12,"color":GRAY},
    {"text":"  文章，要求总结要点","size":12,"color":GRAY},
    {"text":"DeepSeek：精准提取核心","size":12,"color":GRAY},
    {"text":"  观点，结构化呈现 ✓","size":12,"color":GREEN},
], ls=1.2)
pg(s, 7)

# Slide 8
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); decor(s); hdr(s,4,"DeepSeek — 优缺点与使用步骤")
rrect(s, Inches(0.8), Inches(1.2), Inches(3.5), Inches(5.8), CARD, GREEN, Pt(2))
tb(s, Inches(1.0), Inches(1.4), Inches(3), Inches(0.5), "✅ 优点", sz=22, c=GREEN, b=True)
hline(s, Inches(1.0), Inches(1.9), Inches(3), GREEN, Pt(1))
mtb(s, Inches(1.0), Inches(2.1), Inches(3), Inches(4.5),
    [{"text":p,"size":14,"color":WHITE} for p in [
        "• 推理能力行业领先","• 数学/代码能力强",
        "• 免费额度充足","• 支持长上下文（128K）",
        "• 回答逻辑清晰","• 适合深度学习",
        "• 开源生态友好","• 中文理解优秀",
    ]], ls=1.5)

rrect(s, Inches(4.7), Inches(1.2), Inches(3.5), Inches(5.8), CARD, RED, Pt(2))
tb(s, Inches(4.9), Inches(1.4), Inches(3), Inches(0.5), "⚠️ 缺点", sz=22, c=RED, b=True)
hline(s, Inches(4.9), Inches(1.9), Inches(3), RED, Pt(1))
mtb(s, Inches(4.9), Inches(2.1), Inches(3), Inches(4.5),
    [{"text":c,"size":14,"color":WHITE} for c in [
        "• 图像生成能力弱","• 实时信息更新慢",
        "• 创意写作偏理性","• 部分场景响应慢",
        "• 需要一定提问技巧","• 复杂任务可能出错",
        "• 移动端体验一般","• 高级功能需付费",
    ]], ls=1.5)

rrect(s, Inches(8.6), Inches(1.2), Inches(3.9), Inches(5.8), CARD, TEAL, Pt(2))
tb(s, Inches(8.8), Inches(1.4), Inches(3.5), Inches(0.5), "📋 使用步骤", sz=22, c=TEAL, b=True)
hline(s, Inches(8.8), Inches(1.9), Inches(3.5), TEAL, Pt(1))
mtb(s, Inches(8.8), Inches(2.1), Inches(3.5), Inches(4.5), [
    {"text":"Step 1：访问平台","size":15,"color":YELLOW,"bold":True},
    {"text":"  deepseek.com 或下载APP","size":12,"color":GRAY},
    {"text":"","size":6},
    {"text":"Step 2：注册账号","size":15,"color":YELLOW,"bold":True},
    {"text":"  手机号注册即可","size":12,"color":GRAY},
    {"text":"","size":6},
    {"text":"Step 3：选择模型","size":15,"color":YELLOW,"bold":True},
    {"text":"  简单问题→V3模型","size":12,"color":GRAY},
    {"text":"  复杂推理→R1模型","size":12,"color":GRAY},
    {"text":"","size":6},
    {"text":"Step 4：精准提问","size":15,"color":YELLOW,"bold":True},
    {"text":"  给出背景+具体要求","size":12,"color":GRAY},
    {"text":"","size":6},
    {"text":"Step 5：验证结果","size":15,"color":YELLOW,"bold":True},
    {"text":"  重要信息交叉验证","size":12,"color":GRAY},
], ls=1.2)
pg(s, 8)

# ========== SLIDE 9-10: XINGTU ==========
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); decor(s); hdr(s,5,"醒图 — AI智能修图神器")
rrect(s, Inches(0.8), Inches(1.2), Inches(5.8), Inches(5.8), CARD, GREEN, Pt(2))
tb(s, Inches(1.0), Inches(1.4), Inches(5), Inches(0.5), "🎨 醒图", sz=28, c=GREEN, b=True)
hline(s, Inches(1.0), Inches(1.9), Inches(5), GREEN, Pt(1))
mtb(s, Inches(1.0), Inches(2.1), Inches(5.3), Inches(4.5), [
    {"text":"开发商：南京代码科技","size":15,"color":WHITE,"bold":True},
    {"text":"","size":8},
    {"text":"平台：APP（iOS+Android）","size":14,"color":GRAY},
    {"text":"","size":8},
    {"text":"核心功能：","size":15,"color":GREEN,"bold":True},
    {"text":"  • AI智能美颜修图","size":13,"color":GRAY},
    {"text":"  • 一键抠图/换背景","size":13,"color":GRAY},
    {"text":"  • AI风格迁移滤镜","size":13,"color":GRAY},
    {"text":"  • 智能海报设计","size":13,"color":GRAY},
    {"text":"  • AI绘画/文生图","size":13,"color":GRAY},
    {"text":"  • 拼图/排版模板","size":13,"color":GRAY},
    {"text":"  • 批量处理图片","size":13,"color":GRAY},
], ls=1.3)

rrect(s, Inches(7.0), Inches(1.2), Inches(5.5), Inches(5.8), CARD, CYAN, Pt(2))
tb(s, Inches(7.2), Inches(1.4), Inches(5), Inches(0.5), "💡 实用案例", sz=22, c=CYAN, b=True)
hline(s, Inches(7.2), Inches(1.9), Inches(5), CYAN, Pt(1))
mtb(s, Inches(7.2), Inches(2.1), Inches(5), Inches(4.5), [
    {"text":"案例1：商品图片处理","size":15,"color":YELLOW,"bold":True},
    {"text":"电商课作业：给产品图","size":12,"color":GRAY},
    {"text":"  换白色背景+加文字","size":12,"color":GRAY},
    {"text":"醒图：一键抠图→换背景","size":12,"color":GRAY},
    {"text":"  →加文字说明 ✓","size":12,"color":GREEN},
    {"text":"","size":6},
    {"text":"案例2：海报设计","size":15,"color":YELLOW,"bold":True},
    {"text":"学生会活动：制作宣传","size":12,"color":GRAY},
    {"text":"  海报，不会PS怎么办？","size":12,"color":GRAY},
    {"text":"醒图：选模板→换图片→","size":12,"color":GRAY},
    {"text":"  改文字，5分钟搞定 ✓","size":12,"color":GREEN},
    {"text":"","size":6},
    {"text":"案例3：AI艺术照","size":15,"color":YELLOW,"bold":True},
    {"text":"上传自拍→选择风格","size":12,"color":GRAY},
    {"text":"（动漫/油画/素描）","size":12,"color":GRAY},
    {"text":"醒图：AI自动转换风格","size":12,"color":GRAY},
    {"text":"  生成艺术效果 ✓","size":12,"color":GREEN},
], ls=1.2)
pg(s, 9)

# Slide 10
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); decor(s); hdr(s,5,"醒图 — 优缺点与使用步骤")
rrect(s, Inches(0.8), Inches(1.2), Inches(3.5), Inches(5.8), CARD, GREEN, Pt(2))
tb(s, Inches(1.0), Inches(1.4), Inches(3), Inches(0.5), "✅ 优点", sz=22, c=GREEN, b=True)
hline(s, Inches(1.0), Inches(1.9), Inches(3), GREEN, Pt(1))
mtb(s, Inches(1.0), Inches(2.1), Inches(3), Inches(4.5),
    [{"text":p,"size":14,"color":WHITE} for p in [
        "• 操作简单，零基础友好","• 模板丰富，出图快",
        "• AI抠图精准","• 滤镜效果专业",
        "• 免费功能足够多","• 适合手机操作",
        "• 社交分享便捷","• 更新频繁，新功能多",
    ]], ls=1.5)

rrect(s, Inches(4.7), Inches(1.2), Inches(3.5), Inches(5.8), CARD, RED, Pt(2))
tb(s, Inches(4.9), Inches(1.4), Inches(3), Inches(0.5), "⚠️ 缺点", sz=22, c=RED, b=True)
hline(s, Inches(4.9), Inches(1.9), Inches(3), RED, Pt(1))
mtb(s, Inches(4.9), Inches(2.1), Inches(3), Inches(4.5),
    [{"text":c,"size":14,"color":WHITE} for c in [
        "• 高级功能需VIP","• 精细调整不如PS",
        "• 部分模板有水印","• 批量处理能力弱",
        "• 依赖手机性能","• 不适合专业印刷",
        "• 导出画质有压缩","• 广告较多（免费版）",
    ]], ls=1.5)

rrect(s, Inches(8.6), Inches(1.2), Inches(3.9), Inches(5.8), CARD, GREEN, Pt(2))
tb(s, Inches(8.8), Inches(1.4), Inches(3.5), Inches(0.5), "📋 使用步骤", sz=22, c=GREEN, b=True)
hline(s, Inches(8.8), Inches(1.9), Inches(3.5), GREEN, Pt(1))
mtb(s, Inches(8.8), Inches(2.1), Inches(3.5), Inches(4.5), [
    {"text":"Step 1：下载APP","size":15,"color":YELLOW,"bold":True},
    {"text":"  应用商店搜\"醒图\"","size":12,"color":GRAY},
    {"text":"","size":6},
    {"text":"Step 2：导入图片","size":15,"color":YELLOW,"bold":True},
    {"text":"  从相册选择或拍照","size":12,"color":GRAY},
    {"text":"","size":6},
    {"text":"Step 3：选择功能","size":15,"color":YELLOW,"bold":True},
    {"text":"  修图/抠图/海报/模板","size":12,"color":GRAY},
    {"text":"","size":6},
    {"text":"Step 4：调整参数","size":15,"color":YELLOW,"bold":True},
    {"text":"  根据需要调整细节","size":12,"color":GRAY},
    {"text":"","size":6},
    {"text":"Step 5：保存分享","size":15,"color":YELLOW,"bold":True},
    {"text":"  保存到相册或直接","size":12,"color":GRAY},
    {"text":"  分享到社交平台","size":12,"color":GRAY},
], ls=1.2)
pg(s, 10)

# ========== SLIDE 11-12: TONGYI QIANWEN ==========
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); decor(s); hdr(s,6,"通义千问 — 阿里AI助手")
rrect(s, Inches(0.8), Inches(1.2), Inches(5.8), Inches(5.8), CARD, ORANGE, Pt(2))
tb(s, Inches(1.0), Inches(1.4), Inches(5), Inches(0.5), "💬 通义千问", sz=28, c=ORANGE, b=True)
hline(s, Inches(1.0), Inches(1.9), Inches(5), ORANGE, Pt(1))
mtb(s, Inches(1.0), Inches(2.1), Inches(5.3), Inches(4.5), [
    {"text":"开发商：阿里巴巴集团","size":15,"color":WHITE,"bold":True},
    {"text":"","size":8},
    {"text":"平台：网页版 / APP / 钉钉","size":14,"color":GRAY},
    {"text":"网址：tongyi.aliyun.com","size":14,"color":ORANGE},
    {"text":"","size":8},
    {"text":"核心功能：","size":15,"color":ORANGE,"bold":True},
    {"text":"  • 智能对话问答","size":13,"color":GRAY},
    {"text":"  • 长文档分析总结","size":13,"color":GRAY},
    {"text":"  • 代码生成调试","size":13,"color":GRAY},
    {"text":"  • 图像理解（识图）","size":13,"color":GRAY},
    {"text":"  • 文档创作/改写","size":13,"color":GRAY},
    {"text":"  • 数据分析处理","size":13,"color":GRAY},
    {"text":"  • 与钉钉深度整合","size":13,"color":GRAY},
], ls=1.3)

rrect(s, Inches(7.0), Inches(1.2), Inches(5.5), Inches(5.8), CARD, CYAN, Pt(2))
tb(s, Inches(7.2), Inches(1.4), Inches(5), Inches(0.5), "💡 实用案例", sz=22, c=CYAN, b=True)
hline(s, Inches(7.2), Inches(1.9), Inches(5), CYAN, Pt(1))
mtb(s, Inches(7.2), Inches(2.1), Inches(5), Inches(4.5), [
    {"text":"案例1：读书笔记","size":15,"color":YELLOW,"bold":True},
    {"text":"学生：上传课本照片，","size":12,"color":GRAY},
    {"text":"  \"帮我总结这章重点\"","size":12,"color":GRAY},
    {"text":"通义千问：识别文字+","size":12,"color":GRAY},
    {"text":"  提取要点，结构化 ✓","size":12,"color":GREEN},
    {"text":"","size":6},
    {"text":"案例2：PPT大纲生成","size":15,"color":YELLOW,"bold":True},
    {"text":"学生：\"帮我做一份关于","size":12,"color":GRAY},
    {"text":"  电商运营的PPT大纲\"","size":12,"color":GRAY},
    {"text":"通义千问：生成完整大纲","size":12,"color":GRAY},
    {"text":"  含每页标题+要点 ✓","size":12,"color":GREEN},
    {"text":"","size":6},
    {"text":"案例3：英语学习","size":15,"color":YELLOW,"bold":True},
    {"text":"学生：\"帮我分析这篇","size":12,"color":GRAY},
    {"text":"  英语阅读理解的错题\"","size":12,"color":GRAY},
    {"text":"通义千问：逐题解析，","size":12,"color":GRAY},
    {"text":"  讲解语法和词汇 ✓","size":12,"color":GREEN},
], ls=1.2)
pg(s, 11)

# Slide 12
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); decor(s); hdr(s,6,"通义千问 — 优缺点与使用步骤")
rrect(s, Inches(0.8), Inches(1.2), Inches(3.5), Inches(5.8), CARD, GREEN, Pt(2))
tb(s, Inches(1.0), Inches(1.4), Inches(3), Inches(0.5), "✅ 优点", sz=22, c=GREEN, b=True)
hline(s, Inches(1.0), Inches(1.9), Inches(3), GREEN, Pt(1))
mtb(s, Inches(1.0), Inches(2.1), Inches(3), Inches(4.5),
    [{"text":p,"size":14,"color":WHITE} for p in [
        "• 阿里生态，稳定可靠","• 长文本处理能力强",
        "• 图像识别理解好","• 与钉钉/阿里云打通",
        "• 免费额度充足","• 支持文件上传分析",
        "• 中文理解优秀","• 多模态能力强",
    ]], ls=1.5)

rrect(s, Inches(4.7), Inches(1.2), Inches(3.5), Inches(5.8), CARD, RED, Pt(2))
tb(s, Inches(4.9), Inches(1.4), Inches(3), Inches(0.5), "⚠️ 缺点", sz=22, c=RED, b=True)
hline(s, Inches(4.9), Inches(1.9), Inches(3), RED, Pt(1))
mtb(s, Inches(4.9), Inches(2.1), Inches(3), Inches(4.5),
    [{"text":c,"size":14,"color":WHITE} for c in [
        "• 创意写作偏正式","• 部分功能需登录",
        "• 实时信息有延迟","• 数学计算偶有误",
        "• 界面偏商务风","• 移动端功能有限",
        "• 高级模型需排队","• 个性化不足",
    ]], ls=1.5)

rrect(s, Inches(8.6), Inches(1.2), Inches(3.9), Inches(5.8), CARD, ORANGE, Pt(2))
tb(s, Inches(8.8), Inches(1.4), Inches(3.5), Inches(0.5), "📋 使用步骤", sz=22, c=ORANGE, b=True)
hline(s, Inches(8.8), Inches(1.9), Inches(3.5), ORANGE, Pt(1))
mtb(s, Inches(8.8), Inches(2.1), Inches(3.5), Inches(4.5), [
    {"text":"Step 1：访问平台","size":15,"color":YELLOW,"bold":True},
    {"text":"  tongyi.aliyun.com","size":12,"color":GRAY},
    {"text":"","size":6},
    {"text":"Step 2：登录账号","size":15,"color":YELLOW,"bold":True},
    {"text":"  支付宝/淘宝/手机号","size":12,"color":GRAY},
    {"text":"","size":6},
    {"text":"Step 3：选择能力","size":15,"color":YELLOW,"bold":True},
    {"text":"  对话/文档/图片/代码","size":12,"color":GRAY},
    {"text":"","size":6},
    {"text":"Step 4：上传文件","size":15,"color":YELLOW,"bold":True},
    {"text":"  可上传PDF/Word/图片","size":12,"color":GRAY},
    {"text":"","size":6},
    {"text":"Step 5：提问分析","size":15,"color":YELLOW,"bold":True},
    {"text":"  输入具体问题获取","size":12,"color":GRAY},
    {"text":"  AI分析结果","size":12,"color":GRAY},
], ls=1.2)
pg(s, 12)

# ========== SLIDE 13-14: JIAN/YING ==========
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); decor(s); hdr(s,7,"剪映 — AI视频创作工具")
rrect(s, Inches(0.8), Inches(1.2), Inches(5.8), Inches(5.8), CARD, RED, Pt(2))
tb(s, Inches(1.0), Inches(1.4), Inches(5), Inches(0.5), "🎬 剪映", sz=28, c=RED, b=True)
hline(s, Inches(1.0), Inches(1.9), Inches(5), RED, Pt(1))
mtb(s, Inches(1.0), Inches(2.1), Inches(5.3), Inches(4.5), [
    {"text":"开发商：字节跳动","size":15,"color":WHITE,"bold":True},
    {"text":"","size":8},
    {"text":"平台：电脑端 / 手机端","size":14,"color":GRAY},
    {"text":"","size":8},
    {"text":"核心功能：","size":15,"color":RED,"bold":True},
    {"text":"  • AI智能剪辑","size":13,"color":GRAY},
    {"text":"  • 自动字幕生成","size":13,"color":GRAY},
    {"text":"  • AI文案/脚本生成","size":13,"color":GRAY},
    {"text":"  • 智能美颜/滤镜","size":13,"color":GRAY},
    {"text":"  • 图文成片（文字转视频）","size":13,"color":GRAY},
    {"text":"  • 背景音乐智能匹配","size":13,"color":GRAY},
    {"text":"  • 模板一键成片","size":13,"color":GRAY},
], ls=1.3)

rrect(s, Inches(7.0), Inches(1.2), Inches(5.5), Inches(5.8), CARD, CYAN, Pt(2))
tb(s, Inches(7.2), Inches(1.4), Inches(5), Inches(0.5), "💡 实用案例", sz=22, c=CYAN, b=True)
hline(s, Inches(7.2), Inches(1.9), Inches(5), CYAN, Pt(1))
mtb(s, Inches(7.2), Inches(2.1), Inches(5), Inches(4.5), [
    {"text":"案例1：课堂汇报视频","size":15,"color":YELLOW,"bold":True},
    {"text":"作业：制作3分钟产品","size":12,"color":GRAY},
    {"text":"  介绍短视频","size":12,"color":GRAY},
    {"text":"剪映：导入素材→AI自动","size":12,"color":GRAY},
    {"text":"  剪辑→加字幕→导出 ✓","size":12,"color":GREEN},
    {"text":"","size":6},
    {"text":"案例2：图文成片","size":15,"color":YELLOW,"bold":True},
    {"text":"学生：写好产品介绍文案","size":12,"color":GRAY},
    {"text":"剪映：AI自动匹配素材+","size":12,"color":GRAY},
    {"text":"  配音+字幕=完整视频 ✓","size":12,"color":GREEN},
    {"text":"","size":6},
    {"text":"案例3：直播切片","size":15,"color":YELLOW,"bold":True},
    {"text":"电商实训：从1小时直播","size":12,"color":GRAY},
    {"text":"  中剪辑精彩片段","size":12,"color":GRAY},
    {"text":"剪映：AI识别高光时刻+","size":12,"color":GRAY},
    {"text":"  自动剪辑+包装 ✓","size":12,"color":GREEN},
], ls=1.2)
pg(s, 13)

# Slide 14
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); decor(s); hdr(s,7,"剪映 — 优缺点与使用步骤")
rrect(s, Inches(0.8), Inches(1.2), Inches(3.5), Inches(5.8), CARD, GREEN, Pt(2))
tb(s, Inches(1.0), Inches(1.4), Inches(3), Inches(0.5), "✅ 优点", sz=22, c=GREEN, b=True)
hline(s, Inches(1.0), Inches(1.9), Inches(3), GREEN, Pt(1))
mtb(s, Inches(1.0), Inches(2.1), Inches(3), Inches(4.5),
    [{"text":p,"size":14,"color":WHITE} for p in [
        "• 免费使用，功能强大","• 操作简单，门槛低",
        "• AI自动字幕准确率高","• 模板丰富，出片快",
        "• 与抖音无缝对接","• 支持4K导出",
        "• 素材库丰富","• 电脑/手机同步",
    ]], ls=1.5)

rrect(s, Inches(4.7), Inches(1.2), Inches(3.5), Inches(5.8), CARD, RED, Pt(2))
tb(s, Inches(4.9), Inches(1.4), Inches(3), Inches(0.5), "⚠️ 缺点", sz=22, c=RED, b=True)
hline(s, Inches(4.9), Inches(1.9), Inches(3), RED, Pt(1))
mtb(s, Inches(4.9), Inches(2.1), Inches(3), Inches(4.5),
    [{"text":c,"size":14,"color":WHITE} for c in [
        "• 专业级特效有限","• 高级功能需VIP",
        "• 大文件处理较慢","• 自定义程度不如PR",
        "• 部分素材有版权","• 导出有平台水印",
        "• 多轨道编辑不便","• 色彩管理不专业",
    ]], ls=1.5)

rrect(s, Inches(8.6), Inches(1.2), Inches(3.9), Inches(5.8), CARD, RED, Pt(2))
tb(s, Inches(8.8), Inches(1.4), Inches(3.5), Inches(0.5), "📋 使用步骤", sz=22, c=RED, b=True)
hline(s, Inches(8.8), Inches(1.9), Inches(3.5), RED, Pt(1))
mtb(s, Inches(8.8), Inches(2.1), Inches(3.5), Inches(4.5), [
    {"text":"Step 1：下载安装","size":15,"color":YELLOW,"bold":True},
    {"text":"  电脑端/手机端均可","size":12,"color":GRAY},
    {"text":"","size":6},
    {"text":"Step 2：新建项目","size":15,"color":YELLOW,"bold":True},
    {"text":"  导入视频/图片素材","size":12,"color":GRAY},
    {"text":"","size":6},
    {"text":"Step 3：AI剪辑","size":15,"color":YELLOW,"bold":True},
    {"text":"  使用\"图文成片\"或","size":12,"color":GRAY},
    {"text":"  \"智能剪辑\"功能","size":12,"color":GRAY},
    {"text":"","size":6},
    {"text":"Step 4：添加元素","size":15,"color":YELLOW,"bold":True},
    {"text":"  字幕/音乐/特效/转场","size":12,"color":GRAY},
    {"text":"","size":6},
    {"text":"Step 5：导出分享","size":15,"color":YELLOW,"bold":True},
    {"text":"  选择分辨率导出，","size":12,"color":GRAY},
    {"text":"  可直接发抖音","size":12,"color":GRAY},
], ls=1.2)
pg(s, 14)

# ========== SLIDE 15: COMPARISON ==========
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); decor(s); hdr(s,8,"五大AI工具对比总览")
headers = ["功能/工具","📱 豆包","🧠 DeepSeek","🎨 醒图","💬 通义千问","🎬 剪映"]
hcolors = [WHITE,CYAN,TEAL,GREEN,ORANGE,RED]
x = Inches(0.5)
for h, hc in zip(headers, hcolors):
    w = Inches(2.2) if h == "功能/工具" else Inches(2.0)
    rrect(s, x, Inches(1.2), w, Inches(0.5), CARD, hc, Pt(1))
    tb(s, x, Inches(1.25), w, Inches(0.4), h, sz=13, c=hc, b=True, a=PP_ALIGN.CENTER)
    x += w + Inches(0.05)

rows = [
    ["核心定位","AI对话","深度推理","AI修图","多模态AI","AI视频"],
    ["学习辅导","★★★★","★★★★★","★","★★★★","★★"],
    ["创意设计","★★★","★★","★★★★★","★★★","★★★★"],
    ["视频制作","★","★","★★","★★","★★★★★"],
    ["文档处理","★★★★","★★★★★","★","★★★★★","★★"],
    ["上手难度","⭐低","⭐⭐中","⭐低","⭐⭐中","⭐低"],
    ["免费额度","充足","很充足","基础免费","充足","基础免费"],
    ["最佳场景","日常问答","深度学习","图片处理","文件分析","短视频"],
]
y = Inches(1.85)
for row in rows:
    x = Inches(0.5)
    for i, cell in enumerate(row):
        w = Inches(2.2) if i == 0 else Inches(2.0)
        color = WHITE if i == 0 else GRAY
        size = 13 if i == 0 else 12
        rrect(s, x, y, w, Inches(0.55), CARD, DARK, Pt(0.5))
        tb(s, x, y+Inches(0.05), w, Inches(0.45), cell, sz=size, c=color, a=PP_ALIGN.CENTER)
        x += w + Inches(0.05)
    y += Inches(0.6)

rrect(s, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.7), CARD, CYAN, Pt(1))
tb(s, Inches(0.7), Inches(6.55), Inches(12), Inches(0.6),
   "💡 选择建议：学习问答→豆包/通义千问 ｜ 深度学习→DeepSeek ｜ 图片处理→醒图 ｜ 视频制作→剪映",
   sz=14, c=CYAN, a=PP_ALIGN.CENTER)
pg(s, 15)

# ========== SLIDE 16: PROMPT ENGINEERING ==========
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); decor(s); hdr(s,8,"AI工具核心技能：提示词工程")
rrect(s, Inches(0.8), Inches(1.2), Inches(5.8), Inches(2.8), CARD, CYAN, Pt(2))
tb(s, Inches(1.0), Inches(1.4), Inches(5.3), Inches(0.5), "🔑 什么是提示词（Prompt）？", sz=20, c=CYAN, b=True)
hline(s, Inches(1.0), Inches(1.9), Inches(5.3), DIM, Pt(1))
mtb(s, Inches(1.0), Inches(2.1), Inches(5.3), Inches(1.8), [
    "提示词 = 你给AI的「指令」",
    "",
    "好的提示词 → 好的回答",
    "差的提示词 → 差的回答",
    "",
    "👉 提示词工程 = 学会如何向AI提问",
], dsz=15, ls=1.4)

rrect(s, Inches(0.8), Inches(4.2), Inches(5.8), Inches(2.8), CARD, GREEN, Pt(2))
tb(s, Inches(1.0), Inches(4.4), Inches(5.3), Inches(0.5), "📊 好提示词 vs 差提示词", sz=18, c=GREEN, b=True)
hline(s, Inches(1.0), Inches(4.9), Inches(5.3), GREEN, Pt(1))
mtb(s, Inches(1.0), Inches(5.0), Inches(5.3), Inches(1.8), [
    {"text":"❌ 差：\"帮我写作文\"","size":13,"color":RED,"bold":True},
    {"text":"✅ 好：\"写一篇600字记叙文，","size":13,"color":GREEN,"bold":True},
    {"text":"  主题是\"难忘的一天\"，","size":12,"color":GRAY},
    {"text":"  用第一人称，要有细节描写\"","size":12,"color":GRAY},
    {"text":"","size":6},
    {"text":"❌ 差：\"解释光合作用\"","size":13,"color":RED,"bold":True},
    {"text":"✅ 好：\"用初中生能懂的","size":13,"color":GREEN,"bold":True},
    {"text":"  语言解释光合作用，","size":12,"color":GRAY},
    {"text":"  举一个生活中的例子\"","size":12,"color":GRAY},
], ls=1.2)

rrect(s, Inches(7.0), Inches(1.2), Inches(5.5), Inches(5.8), CARD, MAGENTA, Pt(2))
tb(s, Inches(7.2), Inches(1.4), Inches(5), Inches(0.5), "📝 万能提示词公式", sz=22, c=MAGENTA, b=True)
hline(s, Inches(7.2), Inches(1.9), Inches(5), MAGENTA, Pt(1))
mtb(s, Inches(7.2), Inches(2.1), Inches(5), Inches(4.5), [
    {"text":"角色 + 任务 + 要求 + 格式","size":20,"color":YELLOW,"bold":True,"align":PP_ALIGN.CENTER},
    {"text":"","size":10},
    {"text":"📌 角色：你是一位...","size":14,"color":WHITE},
    {"text":"   \"你是一位电商专业老师\"","size":12,"color":GRAY},
    {"text":"","size":8},
    {"text":"📌 任务：请帮我...","size":14,"color":WHITE},
    {"text":"   \"请帮我写一份产品推广文案\"","size":12,"color":GRAY},
    {"text":"","size":8},
    {"text":"📌 要求：字数/风格/重点...","size":14,"color":WHITE},
    {"text":"   \"200字，活泼风格，突出性价比\"","size":12,"color":GRAY},
    {"text":"","size":8},
    {"text":"📌 格式：表格/列表/段落...","size":14,"color":WHITE},
    {"text":"   \"用三点式列出核心卖点\"","size":12,"color":GRAY},
], ls=1.3)
pg(s, 16)

# ========== SLIDE 17: PROMPT PRACTICE ==========
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); decor(s); hdr(s,8,"提示词实战演练")
examples = [
    ("📚 学习场景","你是中职电商专业老师，请帮我列出\"网店运营\"课程的5个核心知识点，每个知识点用1-2句话解释，最后给出学习建议。",CYAN),
    ("✍️ 写作场景","你是一位文案策划，请为一款\"无线蓝牙耳机\"写3条小红书风格的推广文案，每条不超过100字，要活泼有趣，带emoji表情。",GREEN),
    ("🎯 规划场景","我是一名中职高一学生，专业是跨境电商。请帮我制定一个学期学习计划，包含专业技能、文化课、课外活动三个方面，用表格形式呈现。",MAGENTA),
    ("🛒 电商场景","你是一位电商运营专家，请帮我分析\"如何在抖音上推广一款国产护肤品\"，列出5个具体可行的推广策略，每个策略说明预期效果和投入成本。",ORANGE),
]
y = Inches(1.2)
for title, prompt, color in examples:
    rrect(s, Inches(0.8), y, Inches(11.7), Inches(1.3), CARD, color, Pt(1))
    tb(s, Inches(1.0), y+Inches(0.1), Inches(11), Inches(0.4),
       title, sz=18, c=color, b=True)
    tb(s, Inches(1.0), y+Inches(0.5), Inches(11.3), Inches(0.7),
       f"💬 \"{prompt}\"", sz=12, c=GRAY)
    y += Inches(1.45)
pg(s, 17)

# ========== SLIDE 18: PRECAUTIONS ==========
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); decor(s); hdr(s,9,"AI工具使用注意事项 ⚠️")
cards = [
    ("🚫 不要完全依赖AI",["• AI只是辅助工具","• 独立思考最重要","• 不要直接抄AI答案","• 理解比结果更重要","• 培养自己的判断力"],RED,Inches(0.8)),
    ("✅ 要验证AI答案",["• 重要信息交叉验证","• 查资料确认准确性","• AI可能\"幻觉\"编造","• 数学题要自己验算","• 事实性内容查权威源"],YELLOW,Inches(4.3)),
    ("🔒 注意隐私安全",["• 不输入个人敏感信息","• 不上传身份证等证件","• 不泄露家庭住址电话","• 注意账号密码安全","• 了解平台隐私政策"],CYAN,Inches(7.8)),
    ("⚖️ 遵守学术诚信",["• 作业注明AI辅助部分","• 考试不使用AI工具","• 尊重知识产权","• AI生成内容需审核","• 培养诚信品质"],GREEN,Inches(11.3)),
]
for title, items, color, x_pos in cards:
    rrect(s, x_pos, Inches(1.2), Inches(3.2), Inches(5.8), CARD, color, Pt(2))
    tb(s, x_pos+Inches(0.2), Inches(1.4), Inches(2.8), Inches(0.5),
       title, sz=16, c=color, b=True)
    hline(s, x_pos+Inches(0.2), Inches(1.95), Inches(2.8), color, Pt(1))
    mtb(s, x_pos+Inches(0.2), Inches(2.2), Inches(2.8), Inches(4.5),
        [{"text":i,"size":14,"color":WHITE} for i in items], ls=1.6)
pg(s, 18)

# ========== SLIDE 19: PRACTICE TASKS ==========
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); decor(s); hdr(s,10,"课堂实践任务 🎯")
tasks = [
    ("任务一：AI对话体验（15分钟）",CYAN,[
        "📱 打开豆包或通义千问",
        "1. 用\"万能公式\"写一个提示词",
        "2. 向AI提问一个你感兴趣的问题",
        "3. 对比\"好提示词\"和\"差提示词\"的结果差异",
        "4. 记录：AI回答的质量如何？有什么不足？",
        "📝 提交：截图+200字体验报告",
    ]),
    ("任务二：AI修图实战（15分钟）",GREEN,[
        "🎨 打开醒图APP",
        "1. 选择一张自己的照片",
        "2. 使用AI美颜功能调整",
        "3. 使用\"一键抠图\"换背景",
        "4. 添加文字制作成个人海报",
        "📝 提交：原图vs成品图对比",
    ]),
    ("任务三：AI视频创作（15分钟）",RED,[
        "🎬 打开剪映",
        "1. 使用\"图文成片\"功能",
        "2. 输入一段产品介绍文案",
        "3. AI自动生成视频后手动调整",
        "4. 添加字幕和背景音乐",
        "📝 提交：导出视频文件",
    ]),
]
y = Inches(1.2)
for title, color, items in tasks:
    rrect(s, Inches(0.8), y, Inches(11.7), Inches(1.7), CARD, color, Pt(2))
    tb(s, Inches(1.0), y+Inches(0.1), Inches(11), Inches(0.4),
       title, sz=18, c=color, b=True)
    hline(s, Inches(1.0), y+Inches(0.5), Inches(11), color, Pt(1))
    mtb(s, Inches(1.0), y+Inches(0.6), Inches(11), Inches(1.0),
        [{"text":i,"size":14,"color":WHITE} for i in items], ls=1.4)
    y += Inches(1.85)
pg(s, 19)

# ========== SLIDE 20: ROADMAP ==========
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); decor(s); hdr(s,10,"AI工具学习路线图 🗺️")
phases = [
    ("第一阶段\n入门","第1-2周",["• 下载注册5个AI工具","• 体验基本对话功能","• 理解什么是AI","• 学会基础提示词"],CYAN,Inches(0.8)),
    ("第二阶段\n进阶","第3-4周",["• 掌握提示词技巧","• 学习AI修图实操","• 尝试AI视频制作","• 完成课堂实践任务"],GREEN,Inches(3.5)),
    ("第三阶段\n应用","第5-6周",["• AI辅助专业课学习","• 用AI完成课程作业","• 尝试AI创意项目","• 小组协作AI项目"],MAGENTA,Inches(6.2)),
    ("第四阶段\n精通","第7-8周",["• 综合运用多工具","• 独立完成AI作品","• 参加AI技能比赛","• 分享AI使用经验"],ORANGE,Inches(8.9)),
]
for title, time, items, color, x_pos in phases:
    rrect(s, x_pos, Inches(1.2), Inches(2.3), Inches(5.5), CARD, color, Pt(2))
    tb(s, x_pos+Inches(0.1), Inches(1.4), Inches(2.1), Inches(0.7),
       title, sz=18, c=color, b=True, a=PP_ALIGN.CENTER)
    tb(s, x_pos+Inches(0.1), Inches(2.1), Inches(2.1), Inches(0.3),
       time, sz=12, c=GRAY, a=PP_ALIGN.CENTER)
    hline(s, x_pos+Inches(0.2), Inches(2.45), Inches(1.9), color, Pt(1))
    mtb(s, x_pos+Inches(0.15), Inches(2.7), Inches(2.0), Inches(3.5),
        [{"text":i,"size":13,"color":WHITE} for i in items], ls=1.6)
    if x_pos < Inches(8.9):
        tb(s, x_pos+Inches(2.3), Inches(3.5), Inches(0.5), Inches(0.5),
           "→", sz=24, c=color, b=True, a=PP_ALIGN.CENTER)

rrect(s, Inches(0.8), Inches(6.9), Inches(11.7), Inches(0.4), CARD, CYAN, Pt(1))
tb(s, Inches(1.0), Inches(6.92), Inches(11.3), Inches(0.35),
   "💡 学习AI工具的关键：多练、多问、多对比、多总结",
   sz=14, c=CYAN, a=PP_ALIGN.CENTER, b=True)
pg(s, 20)

# ========== SLIDE 21: FAQ ==========
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); decor(s); hdr(s,10,"常见问题 FAQ")
faqs = [
    ("Q1：AI工具都要付费吗？","大部分AI工具有免费版本，足够学生日常使用。高级功能可能需要付费，但学生通常不需要。"),
    ("Q2：用AI写作业算作弊吗？","取决于学校规定。建议：将AI作为学习辅助（如解释概念、提供思路），而非直接抄答案。重要作业应注明AI辅助部分。"),
    ("Q3：AI回答一定准确吗？","不一定！AI可能产生\"幻觉\"（编造信息）。重要信息一定要交叉验证，不要盲目相信。"),
    ("Q4：哪个AI工具最适合学生？","入门推荐豆包（中文好、易上手）；深度学习推荐DeepSeek（推理强）；修图用醒图；视频用剪映。"),
    ("Q5：AI会取代我的工作吗？","AI会取代\"不会用AI的人\"。学会使用AI工具，反而能提升你的竞争力。关键是成为\"会用AI的人\"。"),
]
y = Inches(1.2)
for q, a in faqs:
    rrect(s, Inches(0.8), y, Inches(11.7), Inches(1.0), CARD, DIM, Pt(1))
    tb(s, Inches(1.0), y+Inches(0.05), Inches(11), Inches(0.35), q, sz=15, c=CYAN, b=True)
    tb(s, Inches(1.0), y+Inches(0.4), Inches(11), Inches(0.55), a, sz=13, c=GRAY)
    y += Inches(1.1)
pg(s, 21)

# ========== SLIDE 22: END ==========
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
circ(s, Inches(-2), Inches(-2), Inches(6), RGBColor(0x0D,0x14,0x3A))
circ(s, Inches(10), Inches(4), Inches(5), RGBColor(0x0D,0x14,0x3A))
hline(s, Inches(1.5), Inches(2.5), Inches(10.3), CYAN, Pt(2))
hline(s, Inches(1.5), Inches(5.5), Inches(10.3), MAGENTA, Pt(1))
tb(s, Inches(1.5), Inches(2.8), Inches(10.3), Inches(1.0),
   "谢谢观看！", sz=54, c=WHITE, b=True, a=PP_ALIGN.CENTER)
tb(s, Inches(1.5), Inches(4.0), Inches(10.3), Inches(0.6),
   "AI时代，学会与AI共舞 🕺", sz=26, c=CYAN, a=PP_ALIGN.CENTER)
tb(s, Inches(1.5), Inches(4.8), Inches(10.3), Inches(0.5),
   "掌握AI工具 = 掌握未来竞争力", sz=20, c=MAGENTA, a=PP_ALIGN.CENTER)
tb(s, Inches(1.5), Inches(5.7), Inches(10.3), Inches(0.4),
   "中职高一信息技术课 · AI工具使用指南", sz=16, c=GRAY, a=PP_ALIGN.CENTER)
circ(s, Inches(1.8), Inches(5.8), Inches(0.1), CYAN)
circ(s, Inches(11.4), Inches(5.8), Inches(0.1), MAGENTA)
pg(s, 22)

# ========== ADD SPEAKER NOTES ==========
notes_list = [
    # Slide 1 - Cover
    """各位同学大家好！今天我们一起来学习一门非常实用的课程——AI工具使用指南。
    在开始之前，我想先问大家一个问题：你们平时用手机最多做什么？刷抖音？聊天？打游戏？
    其实，你的手机里已经藏着一个超级助手，只是你可能还没发现。今天，我们就来认识这些AI助手，学会让它们成为你的学习利器。""",

    # Slide 2 - TOC
    """这是我们今天的学习路线图。总共10个部分，从认识AI开始，到5大AI工具的详细讲解，再到实战演练。
    大家可以看到，我们会重点学习5个工具：豆包、DeepSeek、醒图、通义千问、剪映。这些都是目前最主流、最好用的AI工具，而且全部免费。""",

    # Slide 3 - What is AI
    """首先，我们来回答一个问题：什么是AI工具？
    简单来说，AI工具就是由人工智能技术驱动的软件或平台。它们能理解你的语言、识别图像、甚至帮你写作画画。
    为什么中职生要学AI？因为未来90%的工作岗位都会和AI协作。不是AI取代你，而是会用AI的人取代不会用AI的人。
    所以，掌握AI工具，就是掌握未来的竞争力。""",

    # Slide 4 - Overview
    """今天我们要学习的5大AI工具，各有专长。
    豆包擅长日常对话和写作；DeepSeek擅长深度推理和数学编程；醒图是AI修图神器；通义千问是多模态AI助手；剪映是AI视频创作工具。
    它们都是免费的，大家课后都可以下载试用。""",

    # Slide 5 - Doubao intro
    """首先来看豆包，它是字节跳动公司开发的AI助手，就是做抖音的那家公司。
    豆包的功能非常丰富：可以对话问答、写文章、翻译语言、生成图片，还能分析上传的文档。
    我们来看几个实际案例：写作业遇到不会的，可以让豆包用通俗语言解释；作文写不好，可以让豆包帮你润色；毕业了不知道做什么，可以让豆包帮你分析就业方向。""",

    # Slide 6 - Doubao pros/cons/steps
    """豆包的优点很明显：免费、中文理解好、响应快、界面简洁。
    但它也有局限：深度推理能力不如专门的AI，长文本可能遗漏信息，最重要的是——它可能会编造信息，我们叫它"幻觉"。
    所以使用豆包，一定要学会追问和验证。使用步骤很简单：下载APP、注册登录、输入问题、追问优化、复制使用。""",

    # Slide 7 - DeepSeek intro
    """接下来是DeepSeek，中文名"深度求索"，是一家中国AI公司开发的。
    DeepSeek最大的特点是推理能力非常强，特别适合解决数学题、写代码、分析长文档。
    比如数学题，它不仅能给出答案，还会一步步解释每一步的原理，就像老师一样。
    学编程的同学可以用它来写代码，它会生成完整代码并逐行注释。""",

    # Slide 8 - DeepSeek pros/cons/steps
    """DeepSeek的优点：推理能力行业领先、数学代码能力强、免费额度充足、支持超长文本。
    缺点：不擅长图像生成、创意写作偏理性、部分功能需要付费。
    使用技巧：简单问题用V3模型，复杂推理用R1模型。提问时要给出背景和具体要求，得到的答案一定要交叉验证。""",

    # Slide 9 - Xingtu intro
    """醒图是一款手机APP，专门做AI修图的。不会PS的同学，醒图就是你的救星。
    它可以一键抠图换背景、AI美颜、风格迁移、海报设计，甚至AI绘画。
    电商专业的同学经常要做商品图片，用醒图一键抠图换白底，加文字说明，几分钟就搞定。
    做学生会海报？选个模板，换图片改文字，5分钟搞定。""",

    # Slide 10 - Xingtu pros/cons/steps
    """醒图优点：零基础友好、模板丰富、抠图精准、手机就能用。
    缺点：精细调整不如PS、高级功能要VIP、不适合专业印刷。
    使用步骤：下载APP、导入图片、选择功能（修图/抠图/海报）、调整参数、保存分享。""",

    # Slide 11 - Tongyi Qianwen intro
    """通义千问是阿里巴巴的AI助手，和钉钉深度整合，很多功能很强大。
    它最擅长的是长文档分析和图像理解。比如上传课本照片，让它总结这章重点，它能识别文字并提取要点。
    做PPT不知道怎么写大纲？告诉它主题，它帮你生成完整的大纲，每页标题和要点都列好。""",

    # Slide 12 - Tongyi pros/cons/steps
    """通义千问优点：阿里生态稳定可靠、长文本处理强、图像识别好、支持文件上传分析。
    缺点：创意写作偏正式、部分功能需登录、数学计算偶有误。
    使用：tongyi.aliyun.com访问，用支付宝或手机号登录，可以上传PDF、Word、图片让AI分析。""",

    # Slide 13 - Jianying intro
    """剪映也是字节跳动的产品，是目前最流行的AI视频剪辑工具。
    它最厉害的功能是"图文成片"——你写一段文字，AI自动匹配素材、配音、字幕，生成完整视频。
    电商直播的同学，从1小时直播中剪辑精彩片段，剪映AI能自动识别高光时刻并剪辑包装。""",

    # Slide 14 - Jianying pros/cons/steps
    """剪映优点：免费强大、操作简单、AI字幕准确率高、模板丰富、与抖音无缝对接。
    缺点：专业特效有限、大文件处理慢、自定义不如专业软件PR。
    使用：下载安装、新建项目导入素材、使用AI智能剪辑或图文成片、添加字幕音乐特效、导出分享。""",

    # Slide 15 - Comparison
    """这是5大工具的对比总览。大家可以看到，每个工具都有自己的强项。
    学习问答选豆包或通义千问，深度学习选DeepSeek，图片处理选醒图，视频制作选剪映。
    上手难度方面，豆包、醒图、剪映都是低门槛，DeepSeek和通义千问需要一些技巧。
    建议大家根据实际需求选择工具，不是越多越好，而是越适合越好。""",

    # Slide 16 - Prompt Engineering
    """接下来我们学习一个非常重要的技能——提示词工程。
    提示词就是你给AI的指令。好的提示词得到好的回答，差的提示词得到差的回答。
    万能公式是：角色+任务+要求+格式。比如"你是一位电商老师，请帮我写一份产品推广文案，200字，活泼风格，用三点式列出核心卖点"。""",

    # Slide 17 - Prompt Practice
    """我们来看几个实战案例。
    学习场景：让AI扮演老师，列出课程知识点并解释。
    写作场景：让AI扮演文案策划，写小红书风格的推广文案。
    规划场景：让AI帮你制定学期学习计划。
    电商场景：让AI分析抖音推广策略。
    大家课后可以试试这些提示词，感受一下AI的能力。""",

    # Slide 18 - Precautions
    """使用AI工具有四个重要注意事项。
    第一，不要完全依赖AI，独立思考最重要，AI只是辅助工具。
    第二，要验证AI的答案，重要信息一定要交叉验证，因为AI会"幻觉"编造信息。
    第三，注意隐私安全，不要输入个人敏感信息，不要上传身份证等证件。
    第四，遵守学术诚信，作业要注明AI辅助部分，考试不使用AI工具。""",

    # Slide 19 - Practice Tasks
    """接下来是课堂实践环节，有三个任务。
    任务一：AI对话体验，用万能公式写提示词，对比好提示词和差提示词的效果差异。
    任务二：AI修图实战，用醒图做一张个人海报。
    任务三：AI视频创作，用剪映的图文成片功能制作短视频。
    每个任务15分钟，完成后提交作品。""",

    # Slide 20 - Roadmap
    """这是AI工具的学习路线图，分为四个阶段。
    第一阶段入门：下载注册工具，体验基本功能。
    第二阶段进阶：掌握提示词技巧，学习修图和视频。
    第三阶段应用：用AI辅助专业课学习，完成课程作业。
    第四阶段精通：综合运用多工具，参加比赛，分享经验。
    记住：学习AI工具的关键是多练、多问、多对比、多总结。""",

    # Slide 21 - FAQ
    """最后我们回答几个常见问题。
    Q1：AI工具都要付费吗？大部分有免费版本，学生日常使用足够了。
    Q2：用AI写作业算作弊吗？取决于学校规定。建议把AI当学习辅助，不要直接抄答案。
    Q3：AI回答一定准确吗？不一定，可能产生幻觉，重要信息要验证。
    Q4：哪个工具最适合学生？入门选豆包，深度学习选DeepSeek。
    Q5：AI会取代工作吗？AI会取代不会用AI的人，所以要学会用AI。""",

    # Slide 22 - End
    """同学们，今天我们认识了5大AI工具，学习了提示词工程，也做了实战练习。
    记住一句话：AI时代，学会与AI共舞。掌握AI工具，就是掌握未来的竞争力。
    课后请大家下载这5个工具，按照学习路线图开始练习。有问题随时问我。谢谢大家！""",
]

for i, notes_text in enumerate(notes_list):
    if i < len(prs.slides):
        slide = prs.slides[i]
        slide.notes_slide.notes_text_frame.text = notes_text.strip()

# SAVE
out = "/home/admin/.openclaw/workspace/papers/AI工具使用指南_中职高一_带讲稿.pptx"
prs.save(out)
print(f"✅ PPT saved: {out}")
print(f"📄 Total slides: {len(prs.slides)}")
print(f"🎤 Speaker notes added: {len(notes_list)} slides")
