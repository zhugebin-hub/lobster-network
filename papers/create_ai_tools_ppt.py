#!/usr/bin/env python3
"""生成AI工具使用教学PPT（科技风，20+页）"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import math

# ===== 配色方案（科技风） =====
BG_DARK = RGBColor(0x0A, 0x0E, 0x2A)      # 深蓝黑背景
BG_MID = RGBColor(0x11, 0x16, 0x40)        # 中蓝背景
BG_CARD = RGBColor(0x16, 0x1D, 0x50)       # 卡片背景
CYAN = RGBColor(0x00, 0xD4, 0xFF)          # 亮青
CYAN_DIM = RGBColor(0x00, 0x99, 0xBB)      # 暗青
MAGENTA = RGBColor(0xE0, 0x40, 0xFB)       # 品红
GREEN = RGBColor(0x00, 0xE6, 0x76)         # 亮绿
YELLOW = RGBColor(0xFF, 0xEA, 0x00)        # 亮黄
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xDD)
ORANGE = RGBColor(0xFF, 0x91, 0x00)
RED = RGBColor(0xFF, 0x52, 0x52)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW = prs.slide_width
SH = prs.slide_height

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

def hex_to_rgb(h):
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

# ===== 辅助函数 =====

def set_slide_bg(slide, color):
    """设置幻灯片背景色"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, fill_color, border_color=None, border_width=Pt(0)):
    """添加矩形"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.fill.solid()
        shape.line.fill.fore_color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    return shape

def add_rounded_rect(slide, left, top, width, height, fill_color, border_color=None, border_width=Pt(0)):
    """添加圆角矩形"""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.fill.solid()
        shape.line.fill.fore_color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name='微软雅黑'):
    """添加文本框"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_multi_text(slide, left, top, width, height, lines, default_size=16, default_color=WHITE, line_spacing=1.3):
    """添加多行文本"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line_info in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        if isinstance(line_info, str):
            p.text = line_info
            p.font.size = Pt(default_size)
            p.font.color.rgb = default_color
        else:
            p.text = line_info.get('text', '')
            p.font.size = Pt(line_info.get('size', default_size))
            p.font.color.rgb = line_info.get('color', default_color)
            p.font.bold = line_info.get('bold', False)
            p.alignment = line_info.get('align', PP_ALIGN.LEFT)
        p.line_spacing = Pt(int(p.font.size.pt * line_spacing)) if p.font.size else Pt(int(default_size * line_spacing))
    return txBox

def add_circle(slide, left, top, size, fill_color):
    """添加圆形装饰"""
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    # 降低透明度效果通过半透明色实现
    return shape

def add_glow_circle(slide, left, top, size, fill_color, glow_color, glow_size):
    """添加发光圆形装饰"""
    # 外层光晕
    outer = slide.shapes.add_shape(MSO_SHAPE.OVAL, left - glow_size, top - glow_size, size + glow_size*2, size + glow_size*2)
    outer.fill.solid()
    outer.fill.fore_color.rgb = glow_color
    outer.line.fill.background()
    # 内层实心
    inner = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    inner.fill.solid()
    inner.fill.fore_color.rgb = fill_color
    inner.line.fill.background()
    return inner

def add_line(slide, left, top, length, color=CYAN, width_pt=Pt(2)):
    """添加水平线条"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, length, width_pt)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_vline(slide, left, top, length, color=CYAN, width_pt=Pt(2)):
    """添加垂直线条"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width_pt, length)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_decorative_elements(slide):
    """添加科技风装饰元素"""
    # 左上角装饰线
    add_line(slide, Inches(0), Inches(0), Inches(3), CYAN, Pt(3))
    add_line(slide, Inches(0), Inches(0), Inches(2), CYAN, Pt(3))
    # 右下角装饰线
    add_line(slide, Inches(10.333), Inches(7.2), Inches(3), CYAN, Pt(3))
    add_line(slide, Inches(13.033), Inches(5.5), Inches(2), CYAN, Pt(3))
    # 装饰圆
    add_circle(slide, Inches(11.5), Inches(0.3), Inches(0.15), CYAN)
    add_circle(slide, Inches(11.8), Inches(0.5), Inches(0.1), MAGENTA)
    add_circle(slide, Inches(0.3), Inches(6.8), Inches(0.12), CYAN)

def add_slide_number(slide, num, total=22):
    """添加页码"""
    add_text_box(slide, Inches(12.2), Inches(7.0), Inches(1), Inches(0.4),
                 f"{num}/{total}", font_size=12, color=CYAN_DIM, alignment=PP_ALIGN.RIGHT)

def add_section_header(slide, section_num, section_title):
    """添加章节标题（用于内页统一头部）"""
    add_text_box(slide, Inches(0.8), Inches(0.3), Inches(3), Inches(0.5),
                 f"0{section_num}", font_size=24, color=CYAN, bold=True)
    add_text_box(slide, Inches(3.5), Inches(0.3), Inches(8), Inches(0.5),
                 section_title, font_size=20, color=LIGHT_GRAY)
    add_line(slide, Inches(0.8), Inches(0.85), Inches(11.7), CYAN, Pt(1))

def add_card(slide, left, top, width, height, title, content_lines, title_color=CYAN, bg_color=BG_CARD):
    """添加信息卡片"""
    card = add_rounded_rect(slide, left, top, width, height, bg_color, CYAN_DIM, Pt(1))
    add_text_box(slide, left + Inches(0.2), top + Inches(0.15), width - Inches(0.4), Inches(0.4),
                 title, font_size=16, color=title_color, bold=True)
    add_line(slide, left + Inches(0.2), top + Inches(0.55), width - Inches(0.4), CYAN_DIM, Pt(1))
    y = top + Inches(0.7)
    for line in content_lines:
        add_text_box(slide, left + Inches(0.2), y, width - Inches(0.4), Inches(0.35),
                     line, font_size=13, color=LIGHT_GRAY)
        y += Inches(0.32)

# ============================================================
# 第1页：封面
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
set_slide_bg(slide, BG_DARK)

# 背景装饰 - 大圆
add_circle(slide, Inches(-2), Inches(-2), Inches(6), RGBColor(0x0D, 0x14, 0x3A))
add_circle(slide, Inches(10), Inches(4), Inches(5), RGBColor(0x0D, 0x14, 0x3A))
add_circle(slide, Inches(5), Inches(-1), Inches(3), RGBColor(0x0F, 0x18, 0x45))

# 装饰线
add_line(slide, Inches(1.5), Inches(2.8), Inches(10.3), CYAN, Pt(2))
add_line(slide, Inches(1.5), Inches(5.2), Inches(10.3), MAGENTA, Pt(1))

# 主标题
add_text_box(slide, Inches(1.5), Inches(3.0), Inches(10.3), Inches(1.2),
             "AI 工具使用指南", font_size=54, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# 副标题
add_text_box(slide, Inches(1.5), Inches(4.3), Inches(10.3), Inches(0.6),
             "中职高一信息技术课 · 认识身边的AI助手", font_size=24, color=CYAN, alignment=PP_ALIGN.CENTER)

# 底部信息
add_text_box(slide, Inches(1.5), Inches(5.5), Inches(10.3), Inches(0.5),
             "5大AI工具详解 ｜ 实战案例 ｜ 使用步骤 ｜ 优缺点分析",
             font_size=16, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# 装饰小元素
add_circle(slide, Inches(1.8), Inches(5.6), Inches(0.1), CYAN)
add_circle(slide, Inches(11.4), Inches(5.6), Inches(0.1), MAGENTA)

add_slide_number(slide, 1)

# ============================================================
# 第2页：目录
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_decorative_elements(slide)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(5), Inches(0.6),
             "目  录", font_size=36, color=WHITE, bold=True)
add_line(slide, Inches(0.8), Inches(0.9), Inches(4), CYAN, Pt(2))

toc_items = [
    ("01", "什么是AI工具？为什么学？", CYAN),
    ("02", "AI工具全景概览（5大工具）", CYAN),
    ("03", "豆包 — 全能AI对话助手", MAGENTA),
    ("04", "DeepSeek — 深度思考AI", CYAN),
    ("05", "醒图 — AI智能修图神器", GREEN),
    ("06", "通义千问 — 阿里AI助手", CYAN),
    ("07", "剪映 — AI视频创作工具", ORANGE),
    ("08", "工具对比与选择建议", YELLOW),
    ("09", "AI工具使用注意事项", RED),
    ("10", "课堂实践任务", GREEN),
]

y = Inches(1.3)
for num, title, color in toc_items:
    # 编号
    add_text_box(slide, Inches(1.0), y, Inches(0.8), Inches(0.4),
                 num, font_size=20, color=color, bold=True)
    # 标题
    add_text_box(slide, Inches(1.8), y, Inches(8), Inches(0.4),
                 title, font_size=18, color=WHITE)
    # 装饰线
    add_line(slide, Inches(1.0), y + Inches(0.42), Inches(10.5), RGBColor(0x1A, 0x22, 0x55), Pt(1))
    y += Inches(0.55)

add_slide_number(slide, 2)

# ============================================================
# 第3页：什么是AI工具
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_decorative_elements(slide)
add_section_header(slide, 1, "什么是AI工具？")

# 左侧：定义
add_rounded_rect(slide, Inches(0.8), Inches(1.2), Inches(5.5), Inches(5.5), BG_CARD, CYAN_DIM, Pt(1))
add_text_box(slide, Inches(1.0), Inches(1.4), Inches(5), Inches(0.5),
             "🤖 什么是AI工具？", font_size=22, color=CYAN, bold=True)
add_line(slide, Inches(1.0), Inches(1.9), Inches(5), CYAN_DIM, Pt(1))

defs = [
    "AI工具 = 人工智能技术驱动的",
    "软件/平台/应用",
    "",
    "它们能模拟人类的：",
    "  ✓ 语言理解与生成",
    "  ✓ 图像识别与创作",
    "  ✓ 逻辑推理与分析",
    "  ✓ 语音识别与合成",
    "",
    "简单说：AI工具就是",
    "你的「数字助手」！",
]
add_multi_text(slide, Inches(1.0), Inches(2.1), Inches(5), Inches(4),
               defs, default_size=16, line_spacing=1.4)

# 右侧：为什么学
add_rounded_rect(slide, Inches(6.8), Inches(1.2), Inches(5.5), Inches(5.5), BG_CARD, MAGENTA, Pt(1))
add_text_box(slide, Inches(7.0), Inches(1.4), Inches(5), Inches(0.5),
             "🎯 为什么中职生要学AI？", font_size=22, color=MAGENTA, bold=True)
add_line(slide, Inches(7.0), Inches(1.9), Inches(5), RGBColor(0xAA, 0x30, 0xCC), Pt(1))

whys = [
    "1. 未来职场必备技能",
    "   → 90%岗位将与AI协作",
    "",
    "2. 提升学习效率",
    "   → 快速查找资料、解答疑问",
    "",
    "3. 激发创意思维",
    "   → AI帮你突破灵感瓶颈",
    "",
    "4. 增强就业竞争力",
    "   → 掌握AI=掌握未来",
]
add_multi_text(slide, Inches(7.0), Inches(2.1), Inches(5), Inches(4),
               whys, default_size=15, line_spacing=1.3)

add_slide_number(slide, 3)

# ============================================================
# 第4页：AI工具全景概览
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_decorative_elements(slide)
add_section_header(slide, 2, "AI工具全景概览")

tools_overview = [
    ("📱 豆包", "字节跳动", "AI对话/写作/问答", CYAN),
    ("🧠 DeepSeek", "深度求索", "深度推理/代码/分析", RGBColor(0x64, 0xFF, 0xDA)),
    ("🎨 醒图", "南京代码", "AI修图/设计/美颜", GREEN),
    ("💬 通义千问", "阿里巴巴", "对话/创作/多模态", ORANGE),
    ("🎬 剪映", "字节跳动", "AI视频剪辑/特效", RED),
]

x = Inches(0.6)
for i, (name, company, desc, color) in enumerate(tools_overview):
    card = add_rounded_rect(slide, x, Inches(1.3), Inches(2.3), Inches(5.2), BG_CARD, color, Pt(2))
    # 图标区域
    add_circle(slide, x + Inches(0.7), Inches(1.6), Inches(0.9), color)
    add_text_box(slide, x + Inches(0.15), Inches(2.6), Inches(2.0), Inches(0.8),
                 name, font_size=22, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.15), Inches(3.5), Inches(2.0), Inches(0.4),
                 f"开发商：{company}", font_size=12, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
    add_line(slide, x + Inches(0.3), Inches(4.0), Inches(1.7), color, Pt(1))
    add_text_box(slide, x + Inches(0.15), Inches(4.2), Inches(2.0), Inches(0.4),
                 "核心功能：", font_size=12, color=CYAN_DIM, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.15), Inches(4.5), Inches(2.0), Inches(0.8),
                 desc, font_size=14, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.15), Inches(5.5), Inches(2.0), Inches(0.4),
                 "免费使用 ✓", font_size=12, color=GREEN, alignment=PP_ALIGN.CENTER)
    x += Inches(2.5)

add_slide_number(slide, 4)

# ============================================================
# 第5-6页：豆包
# ============================================================
# 第5页：豆包介绍+案例
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_decorative_elements(slide)
add_section_header(slide, 3, "豆包 — 字节跳动AI助手")

# 左侧：基本信息
add_rounded_rect(slide, Inches(0.8), Inches(1.2), Inches(5.8), Inches(5.8), BG_CARD, CYAN, Pt(2))
add_text_box(slide, Inches(1.0), Inches(1.4), Inches(5), Inches(0.5),
             "📱 豆包", font_size=28, color=CYAN, bold=True)
add_line(slide, Inches(1.0), Inches(1.95), Inches(5), CYAN, Pt(1))

baobao_info = [
    {"text": "开发商：字节跳动（抖音母公司）", "size": 15, "color": WHITE, "bold": True},
    {"text": "", "size": 8},
    {"text": "平台：网页版 / APP（iOS+Android）", "size": 14, "color": LIGHT_GRAY},
    {"text": "网址：www.doubao.com", "size": 14, "color": CYAN},
    {"text": "", "size": 8},
    {"text": "核心功能：", "size": 15, "color": CYAN, "bold": True},
    {"text": "  • 智能对话问答", "size": 13, "color": LIGHT_GRAY},
    {"text": "  • 文章写作/改写", "size": 13, "color": LIGHT_GRAY},
    {"text": "  • 知识讲解/ tutoring", "size": 13, "color": LIGHT_GRAY},
    {"text": "  • 图片生成（文生图）", "size": 13, "color": LIGHT_GRAY},
    {"text": "  • 语言翻译", "size": 13, "color": LIGHT_GRAY},
    {"text": "  • 文件分析（上传文档）", "size": 13, "color": LIGHT_GRAY},
]
add_multi_text(slide, Inches(1.0), Inches(2.1), Inches(5.3), Inches(4.5),
               baobao_info, line_spacing=1.3)

# 右侧：使用案例
add_rounded_rect(slide, Inches(7.0), Inches(1.2), Inches(5.5), Inches(5.8), BG_CARD, GREEN, Pt(2))
add_text_box(slide, Inches(7.2), Inches(1.4), Inches(5), Inches(0.5),
             "💡 实用案例", font_size=22, color=GREEN, bold=True)
add_line(slide, Inches(7.2), Inches(1.95), Inches(5), GREEN, Pt(1))

cases = [
    {"text": "案例1：写作业助手", "size": 15, "color": YELLOW, "bold": True},
    {"text": "学生：\"帮我解释光合作", "size": 12, "color": LIGHT_GRAY},
    {"text": "用的过程，用通俗语言\"", "size": 12, "color": LIGHT_GRAY},
    {"text": "豆包：生成清晰易懂的", "size": 12, "color": LIGHT_GRAY},
    {"text": "  分步骤解释 ✓", "size": 12, "color": GREEN},
    {"text": "", "size": 6},
    {"text": "案例2：作文润色", "size": 15, "color": YELLOW, "bold": True},
    {"text": "学生：把这段文字改写", "size": 12, "color": LIGHT_GRAY},
    {"text": "得更生动，加入修辞手法", "size": 12, "color": LIGHT_GRAY},
    {"text": "豆包：自动优化表达，", "size": 12, "color": LIGHT_GRAY},
    {"text": "  提升文采 ✓", "size": 12, "color": GREEN},
    {"text": "", "size": 6},
    {"text": "案例3：职业规划", "size": 15, "color": YELLOW, "bold": True},
    {"text": "学生：\"中职电商专业", "size": 12, "color": LIGHT_GRAY},
    {"text": " 毕业后可以做什么？\"", "size": 12, "color": LIGHT_GRAY},
    {"text": "豆包：列出就业方向+技能", "size": 12, "color": LIGHT_GRAY},
    {"text": "  建议+薪资参考 ✓", "size": 12, "color": GREEN},
]
add_multi_text(slide, Inches(7.2), Inches(2.1), Inches(5), Inches(4.5),
               cases, line_spacing=1.2)

add_slide_number(slide, 5)

# 第6页：豆包优缺点+使用步骤
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_decorative_elements(slide)
add_section_header(slide, 3, "豆包 — 优缺点与使用步骤")

# 左侧：优点
add_rounded_rect(slide, Inches(0.8), Inches(1.2), Inches(3.5), Inches(5.8), BG_CARD, GREEN, Pt(2))
add_text_box(slide, Inches(1.0), Inches(1.4), Inches(3), Inches(0.5),
             "✅ 优点", font_size=22, color=GREEN, bold=True)
add_line(slide, Inches(1.0), Inches(1.9), Inches(3), GREEN, Pt(1))
pros = [
    "• 免费使用，无门槛",
    "• 中文理解能力强",
    "• 响应速度快",
    "• 支持多种格式输入",
    "• 与抖音生态联动",
    "• 适合日常学习问答",
    "• 界面简洁易用",
    "• 支持图片生成",
]
add_multi_text(slide, Inches(1.0), Inches(2.1), Inches(3), Inches(4.5),
               [{"text": p, "size": 14, "color": WHITE} for p in pros], line_spacing=1.5)

# 中间：缺点
add_rounded_rect(slide, Inches(4.7), Inches(1.2), Inches(3.5), Inches(5.8), BG_CARD, RED, Pt(2))
add_text_box(slide, Inches(4.9), Inches(1.4), Inches(3), Inches(0.5),
             "⚠️ 缺点", font_size=22, color=RED, bold=True)
add_line(slide, Inches(4.9), Inches(1.9), Inches(3), RED, Pt(1))
cons = [
    "• 深度推理能力有限",
    "• 长文本可能遗漏",
    "• 专业领域知识不足",
    "• 可能产生\"幻觉\"",
    "  （编造信息）",
    "• 不能替代独立思考",
    "• 免费额度有限制",
    "• 依赖网络连接",
]
add_multi_text(slide, Inches(4.9), Inches(2.1), Inches(3), Inches(4.5),
               [{"text": c, "size": 14, "color": WHITE} for c in cons], line_spacing=1.5)

# 右侧：使用步骤
add_rounded_rect(slide, Inches(8.6), Inches(1.2), Inches(3.9), Inches(5.8), BG_CARD, CYAN, Pt(2))
add_text_box(slide, Inches(8.8), Inches(1.4), Inches(3.5), Inches(0.5),
             "📋 使用步骤", font_size=22, color=CYAN, bold=True)
add_line(slide, Inches(8.8), Inches(1.9), Inches(3.5), CYAN, Pt(1))

steps = [
    {"text": "Step 1：下载/访问", "size": 15, "color": YELLOW, "bold": True},
    {"text": "  手机应用商店搜\"豆包\"", "size": 12, "color": LIGHT_GRAY},
    {"text": "  或网页访问 doubao.com", "size": 12, "color": LIGHT_GRAY},
    {"text": "", "size": 6},
    {"text": "Step 2：注册登录", "size": 15, "color": YELLOW, "bold": True},
    {"text": "  手机号/抖音账号登录", "size": 12, "color": LIGHT_GRAY},
    {"text": "", "size": 6},
    {"text": "Step 3：输入问题", "size": 15, "color": YELLOW, "bold": True},
    {"text": "  在对话框输入你的问题", "size": 12, "color": LIGHT_GRAY},
    {"text": "  问题越具体，回答越好", "size": 12, "color": LIGHT_GRAY},
    {"text": "", "size": 6},
    {"text": "Step 4：追问优化", "size": 15, "color": YELLOW, "bold": True},
    {"text": "  对回答不满意可继续追问", "size": 12, "color": LIGHT_GRAY},
    {"text": "", "size": 6},
    {"text": "Step 5：复制使用", "size": 15, "color": YELLOW, "bold": True},
    {"text": "  复制有用内容到自己的", "size": 12, "color": LIGHT_GRAY},
    {"text": "  文档中整理使用", "size": 12, "color": LIGHT_GRAY},
]
add_multi_text(slide, Inches(8.8), Inches(2.1), Inches(3.5), Inches(4.5),
               steps, line_spacing=1.2)

add_slide_number(slide, 6)

# ============================================================
# 第7-8页：DeepSeek
# ============================================================
# 第7页：DeepSeek介绍+案例
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_decorative_elements(slide)
add_section_header(slide, 4, "DeepSeek — 深度推理AI")

# 左侧：基本信息
add_rounded_rect(slide, Inches(0.8), Inches(1.2), Inches(5.8), Inches(5.8), BG_CARD, RGBColor(0x64, 0xFF, 0xDA), Pt(2))
add_text_box(slide, Inches(1.0), Inches(1.4), Inches(5), Inches(0.5),
             "🧠 DeepSeek", font_size=28, color=RGBColor(0x64, 0xFF, 0xDA), bold=True)
add_line(slide, Inches(1.0), Inches(1.95), Inches(5), RGBColor(0x64, 0xFF, 0xDA), Pt(1))

ds_info = [
    {"text": "开发商：深度求索（中国AI公司）", "size": 15, "color": WHITE, "bold": True},
    {"text": "", "size": 8},
    {"text": "平台：网页版 / APP", "size": 14, "color": LIGHT_GRAY},
    {"text": "网址：www.deepseek.com", "size": 14, "color": RGBColor(0x64, 0xFF, 0xDA)},
    {"text": "", "size": 8},
    {"text": "核心特色：", "size": 15, "color": RGBColor(0x64, 0xFF, 0xDA), "bold": True},
    {"text": "  • 深度推理能力强大", "size": 13, "color": LIGHT_GRAY},
    {"text": "  • 数学/逻辑题擅长", "size": 13, "color": LIGHT_GRAY},
    {"text": "  • 代码编写与调试", "size": 13, "color": LIGHT_GRAY},
    {"text": "  • 长文本分析处理", "size": 13, "color": LIGHT_GRAY},
    {"text": "  • 多语言翻译精准", "size": 13, "color": LIGHT_GRAY},
    {"text": "  • 免费额度 generous", "size": 13, "color": LIGHT_GRAY},
]
add_multi_text(slide, Inches(1.0), Inches(2.1), Inches(5.3), Inches(4.5),
               ds_info, line_spacing=1.3)

# 右侧：使用案例
add_rounded_rect(slide, Inches(7.0), Inches(1.2), Inches(5.5), Inches(5.8), BG_CARD, CYAN, Pt(2))
add_text_box(slide, Inches(7.2), Inches(1.4), Inches(5), Inches(0.5),
             "💡 实用案例", font_size=22, color=CYAN, bold=True)
add_line(slide, Inches(7.2), Inches(1.95), Inches(5), CYAN, Pt(1))

ds_cases = [
    {"text": "案例1：数学题详解", "size": 15, "color": YELLOW, "bold": True},
    {"text": "学生：\"解方程 2x²+5x-3=0", "size": 12, "color": LIGHT_GRAY},
    {"text": "  请给出详细步骤\"", "size": 12, "color": LIGHT_GRAY},
    {"text": "DeepSeek：逐步推导，", "size": 12, "color": LIGHT_GRAY},
    {"text": "  解释每一步原理 ✓", "size": 12, "color": GREEN},
    {"text": "", "size": 6},
    {"text": "案例2：编程入门", "size": 15, "color": YELLOW, "bold": True},
    {"text": "学生：\"用Python写一个", "size": 12, "color": LIGHT_GRAY},
    {"text": "  猜数字小游戏\"", "size": 12, "color": LIGHT_GRAY},
    {"text": "DeepSeek：生成完整代码+", "size": 12, "color": LIGHT_GRAY},
    {"text": "  逐行注释讲解 ✓", "size": 12, "color": GREEN},
    {"text": "", "size": 6},
    {"text": "案例3：长文档分析", "size": 15, "color": YELLOW, "bold": True},
    {"text": "学生：上传一篇3000字", "size": 12, "color": LIGHT_GRAY},
    {"text": "  文章，要求总结要点", "size": 12, "color": LIGHT_GRAY},
    {"text": "DeepSeek：精准提取核心", "size": 12, "color": LIGHT_GRAY},
    {"text": "  观点，结构化呈现 ✓", "size": 12, "color": GREEN},
]
add_multi_text(slide, Inches(7.2), Inches(2.1), Inches(5), Inches(4.5),
               ds_cases, line_spacing=1.2)

add_slide_number(slide, 7)

# 第8页：DeepSeek优缺点+使用步骤
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_decorative_elements(slide)
add_section_header(slide, 4, "DeepSeek — 优缺点与使用步骤")

# 三栏布局
add_rounded_rect(slide, Inches(0.8), Inches(1.2), Inches(3.5), Inches(5.8), BG_CARD, GREEN, Pt(2))
add_text_box(slide, Inches(1.0), Inches(1.4), Inches(3), Inches(0.5),
             "✅ 优点", font_size=22, color=GREEN, bold=True)
add_line(slide, Inches(1.0), Inches(1.9), Inches(3), GREEN, Pt(1))
ds_pros = [
    "• 推理能力行业领先",
    "• 数学/代码能力强",
    "• 免费额度 generous",
    "• 支持长上下文（128K）",
    "• 回答逻辑清晰",
    "• 适合深度学习",
    "• 开源生态友好",
    "• 中文理解优秀",
]
add_multi_text(slide, Inches(1.0), Inches(2.1), Inches(3), Inches(4.5),
               [{"text": p, "size": 14, "color": WHITE} for p in ds_pros], line_spacing=1.5)

add_rounded_rect(slide, Inches(4.7), Inches(1.2), Inches(3.5), Inches(5.8), BG_CARD, RED, Pt(2))
add_text_box(slide, Inches(4.9), Inches(1.4), Inches(3), Inches(0.5),
             "⚠️ 缺点", font_size=22, color=RED, bold=True)
add_line(slide, Inches(4.9), Inches(1.9), Inches(3), RED, Pt(1))
ds_cons = [
    "• 图像生成能力弱",
    "• 实时信息更新慢",
    "• 创意写作偏理性",
    "• 部分场景响应慢",
    "• 需要一定提问技巧",
    "• 复杂任务可能出错",
    "• 移动端体验一般",
    "• 高级功能需付费",
]
add_multi_text(slide, Inches(4.9), Inches(2.1), Inches(3), Inches(4.5),
               [{"text": c, "size": 14, "color": WHITE} for c in ds_cons], line_spacing=1.5)

add_rounded_rect(slide, Inches(8.6), Inches(1.2), Inches(3.9), Inches(5.8), BG_CARD, RGBColor(0x64, 0xFF, 0xDA), Pt(2))
add_text_box(slide, Inches(8.8), Inches(1.4), Inches(3.5), Inches(0.5),
             "📋 使用步骤", font_size=22, color=RGBColor(0x64, 0xFF, 0xDA), bold=True)
add_line(slide, Inches(8.8), Inches(1.9), Inches(3.5), RGBColor(0x64, 0xFF, 0xDA), Pt(1))
ds_steps = [
    {"text": "Step 1：访问平台", "size": 15, "color": YELLOW, "bold": True},
    {"text": "  deepseek.com 或下载APP", "size": 12, "color": LIGHT_GRAY},
    {"text": "", "size": 6},
    {"text": "Step 2：注册账号", "size": 15, "color": YELLOW, "bold": True},
    {"text": "  手机号注册即可", "size": 12, "color": LIGHT_GRAY},
    {"text": "", "size": 6},
    {"text": "Step 3：选择模型", "size": 15, "color": YELLOW, "bold": True},
    {"text": "  简单问题→V3模型", "size": 12, "color": LIGHT_GRAY},
    {"text": "  复杂推理→R1模型", "size": 12, "color": LIGHT_GRAY},
    {"text": "", "size": 6},
    {"text": "Step 4：精准提问", "size": 15, "color": YELLOW, "bold": True},
    {"text": "  给出背景+具体要求", "size": 12, "color": LIGHT_GRAY},
    {"text": "", "size": 6},
    {"text": "Step 5：验证结果", "size": 15, "color": YELLOW, "bold": True},
    {"text": "  重要信息交叉验证", "size": 12, "color": LIGHT_GRAY},
]
add_multi_text(slide, Inches(8.8), Inches(2.1), Inches(3.5), Inches(4.5),
               ds_steps, line_spacing=1.2)

add_slide_number(slide, 8)

# ============================================================
# 第9-10页：醒图
# ============================================================
# 第9页：醒图介绍+案例
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_decorative_elements(slide)
add_section_header(slide, 5, "醒图 — AI智能修图神器")

add_rounded_rect(slide, Inches(0.8), Inches(1.2), Inches(5.8), Inches(5.8), BG_CARD, GREEN, Pt(2))
add_text_box(slide, Inches(1.0), Inches(1.4), Inches(5), Inches(0.5),
             "🎨 醒图", font_size=28, color=GREEN, bold=True)
add_line(slide, Inches(1.0), Inches(1.95), Inches(5), GREEN, Pt(1))

xt_info = [
    {"text": "开发商：南京代码科技", "size": 15, "color": WHITE, "bold": True},
    {"text": "", "size": 8},
    {"text": "平台：APP（iOS+Android）", "size": 14, "color": LIGHT_GRAY},
    {"text": "", "size": 8},
    {"text": "核心功能：", "size": 15, "color": GREEN, "bold": True},
    {"text": "  • AI智能美颜修图", "size": 13, "color": LIGHT_GRAY},
    {"text": "  • 一键抠图/换背景", "size": 13, "color": LIGHT_GRAY},
    {"text": "  • AI风格迁移滤镜", "size": 13, "color": LIGHT_GRAY},
    {"text": "  • 智能海报设计", "size": 13, "color": LIGHT_GRAY},
    {"text": "  • AI绘画/文生图", "size": 13, "color": LIGHT_GRAY},
    {"text": "  • 拼图/排版模板", "size": 13, "color": LIGHT_GRAY},
    {"text": "  • 批量处理图片", "size": 13, "color": LIGHT_GRAY},
]
add_multi_text(slide, Inches(1.0), Inches(2.1), Inches(5.3), Inches(4.5),
               xt_info, line_spacing=1.3)

add_rounded_rect(slide, Inches(7.0), Inches(1.2), Inches(5.5), Inches(5.8), BG_CARD, CYAN, Pt(2))
add_text_box(slide, Inches(7.2), Inches(1.4), Inches(5), Inches(0.5),
             "💡 实用案例", font_size=22, color=CYAN, bold=True)
add_line(slide, Inches(7.2), Inches(1.95), Inches(5), CYAN, Pt(1))

xt_cases = [
    {"text": "案例1：商品图片处理", "size": 15, "color": YELLOW, "bold": True},
    {"text": "电商课作业：给产品图", "size": 12, "color": LIGHT_GRAY},
    {"text": "  换白色背景+加文字", "size": 12, "color": LIGHT_GRAY},
    {"text": "醒图：一键抠图→换背景", "size": 12, "color": LIGHT_GRAY},
    {"text": "  →加文字说明 ✓", "size": 12, "color": GREEN},
    {"text": "", "size": 6},
    {"text": "案例2：海报设计", "size": 15, "color": YELLOW, "bold": True},
    {"text": "学生会活动：制作宣传", "size": 12, "color": LIGHT_GRAY},
    {"text": "  海报，不会PS怎么办？", "size": 12, "color": LIGHT_GRAY},
    {"text": "醒图：选模板→换图片→", "size": 12, "color": LIGHT_GRAY},
    {"text": "  改文字，5分钟搞定 ✓", "size": 12, "color": GREEN},
    {"text": "", "size": 6},
    {"text": "案例3：AI艺术照", "size": 15, "color": YELLOW, "bold": True},
    {"text": "上传自拍→选择风格", "size": 12, "color": LIGHT_GRAY},
    {"text": "（动漫/油画/素描）", "size": 12, "color": LIGHT_GRAY},
    {"text": "醒图：AI自动转换风格", "size": 12, "color": LIGHT_GRAY},
    {"text": "  生成艺术效果 ✓", "size": 12, "color": GREEN},
]
add_multi_text(slide, Inches(7.2), Inches(2.1), Inches(5), Inches(4.5),
               xt_cases, line_spacing=1.2)

add_slide_number(slide, 9)

# 第10页：醒图优缺点+使用步骤
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_decorative_elements(slide)
add_section_header(slide, 5, "醒图 — 优缺点与使用步骤")

add_rounded_rect(slide, Inches(0.8), Inches(1.2), Inches(3.5), Inches(5.8), BG_CARD, GREEN, Pt(2))
add_text_box(slide, Inches(1.0), Inches(1.4), Inches(3), Inches(0.5),
             "✅ 优点", font_size=22, color=GREEN, bold=True)
add_line(slide, Inches(1.0), Inches(1.9), Inches(3), GREEN, Pt(1))
xt_pros = [
    "• 操作简单，零基础友好",
    "• 模板丰富，出图快",
    "• AI抠图精准",
    "• 滤镜效果专业",
    "• 免费功能足够多",
    "• 适合手机操作",
    "• 社交分享便捷",
    "• 更新频繁，新功能多",
]
add_multi_text(slide, Inches(1.0), Inches(2.1), Inches(3), Inches(4.5),
               [{"text": p, "size": 14, "color": WHITE} for p in xt_pros], line_spacing=1.5)

add_rounded_rect(slide, Inches(4.7), Inches(1.2), Inches(3.5), Inches(5.8), BG_CARD, RED, Pt(2))
add_text_box(slide, Inches(4.9), Inches(1.4), Inches(3), Inches(0.5),
             "⚠️ 缺点", font_size=22, color=RED, bold=True)
add_line(slide, Inches(4.9), Inches(1.9), Inches(3), RED, Pt(1))
xt_cons = [
    "• 高级功能需VIP",
    "• 精细调整不如PS",
    "• 部分模板有水印",
    "• 批量处理能力弱",
    "• 依赖手机性能",
    "• 不适合专业印刷",
    "• 导出画质有压缩",
    "• 广告较多（免费版）",
]
add_multi_text(slide, Inches(4.9), Inches(2.1), Inches(3), Inches(4.5),
               [{"text": c, "size": 14, "color": WHITE} for c in xt_cons], line_spacing=1.5)

add_rounded_rect(slide, Inches(8.6), Inches(1.2), Inches(3.9), Inches(5.8), BG_CARD, GREEN, Pt(2))
add_text_box(slide, Inches(8.8), Inches(1.4), Inches(3.5), Inches(0.5),
             "📋 使用步骤", font_size=22, color=GREEN, bold=True)
add_line(slide, Inches(8.8), Inches(1.9), Inches(3.5), GREEN, Pt(1))
xt_steps = [
    {"text": "Step 1：下载APP", "size": 15, "color": YELLOW, "bold": True},
    {"text": "  应用商店搜\"醒图\"", "size": 12, "color": LIGHT_GRAY},
    {"text": "", "size": 6},
    {"text": "Step 2：导入图片", "size": 15, "color": YELLOW, "bold": True},
    {"text": "  从相册选择或拍照", "size": 12, "color": LIGHT_GRAY},
    {"text": "", "size": 6},
    {"text": "Step 3：选择功能", "size": 15, "color": YELLOW, "bold": True},
    {"text": "  修图/抠图/海报/模板", "size": 12, "color": LIGHT_GRAY},
    {"text": "", "size": 6},
    {"text": "Step 4：调整参数", "size": 15, "color": YELLOW, "bold": True},
    {"text": "  根据需要调整细节", "size": 12, "color": LIGHT_GRAY},
    {"text": "", "size": 6},
    {"text": "Step 5：保存分享", "size": 15, "color": YELLOW, "bold": True},
    {"text": "  保存到相册或直接", "size": 12, "color": LIGHT_GRAY},
    {"text": "  分享到社交平台", "size": 12, "color": LIGHT_GRAY},
]
add_multi_text(slide, Inches(8.8), Inches(2.1), Inches(3.5), Inches(4.5),
               xt_steps, line_spacing=1.2)

add_slide_number(slide, 10)

# ============================================================
# 第11-12页：通义千问
# ============================================================
# 第11页：通义千问介绍+案例
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_decorative_elements(slide)
add_section_header(slide, 6, "通义千问 — 阿里AI助手")

add_rounded_rect(slide, Inches(0.8), Inches(1.2), Inches(5.8), Inches(5.8), BG_CARD, ORANGE, Pt(2))
add_text_box(slide, Inches(1.0), Inches(1.4), Inches(5), Inches(0.5),
             "💬 通义千问", font_size=28, color=ORANGE, bold=True)
add_line(slide, Inches(1.0), Inches(1.95), Inches(5), ORANGE, Pt(1))

qwen_info = [
    {"text": "开发商：阿里巴巴集团", "size": 15, "color": WHITE, "bold": True},
    {"text": "", "size": 8},
    {"text": "平台：网页版 / APP / 钉钉", "size": 14, "color": LIGHT_GRAY},
    {"text": "网址：tongyi.aliyun.com", "size": 14, "color": ORANGE},
    {"text": "", "size": 8},
    {"text": "核心功能：", "size": 15, "color": ORANGE, "bold": True},
    {"text": "  • 智能对话问答", "size": 13, "color": LIGHT_GRAY},
    {"text": "  • 长文档分析总结", "size": 13, "color": LIGHT_GRAY},
    {"text": "  • 代码生成调试", "size": 13, "color": LIGHT_GRAY},
    {"text": "  • 图像理解（识图）", "size": 13, "color": LIGHT_GRAY},
    {"text": "  • 文档创作/改写", "size": 13, "color": LIGHT_GRAY},
    {"text": "  • 数据分析处理", "size": 13, "color": LIGHT_GRAY},
    {"text": "  • 与钉钉深度整合", "size": 13, "color": LIGHT_GRAY},
]
add_multi_text(slide, Inches(1.0), Inches(2.1), Inches(5.3), Inches(4.5),
               qwen_info, line_spacing=1.3)

add_rounded_rect(slide, Inches(7.0), Inches(1.2), Inches(5.5), Inches(5.8), BG_CARD, CYAN, Pt(2))
add_text_box(slide, Inches(7.2), Inches(1.4), Inches(5), Inches(0.5),
             "💡 实用案例", font_size=22, color=CYAN, bold=True)
add_line(slide, Inches(7.2), Inches(1.95), Inches(5), CYAN, Pt(1))

qwen_cases = [
    {"text": "案例1：读书笔记", "size": 15, "color": YELLOW, "bold": True},
    {"text": "学生：上传课本照片，", "size": 12, "color": LIGHT_GRAY},
    {"text": "  \"帮我总结这章重点\"", "size": 12, "color": LIGHT_GRAY},
    {"text": "通义千问：识别文字+", "size": 12, "color": LIGHT_GRAY},
    {"text": "  提取要点，结构化 ✓", "size": 12, "color": GREEN},
    {"text": "", "size": 6},
    {"text": "案例2：PPT大纲生成", "size": 15, "color": YELLOW, "bold": True},
    {"text": "学生：\"帮我做一份关于", "size": 12, "color": LIGHT_GRAY},
    {"text": "  电商运营的PPT大纲\"", "size": 12, "color": LIGHT_GRAY},
    {"text": "通义千问：生成完整大纲", "size": 12, "color": LIGHT_GRAY},
    {"text": "  含每页标题+要点 ✓", "size": 12, "color": GREEN},
    {"text": "", "size": 6},
    {"text": "案例3：英语学习", "size": 15, "color": YELLOW, "bold": True},
    {"text": "学生：\"帮我分析这篇", "size": 12, "color": LIGHT_GRAY},
    {"text": "  英语阅读理解的错题\"", "size": 12, "color": LIGHT_GRAY},
    {"text": "通义千问：逐题解析，", "size": 12, "color": LIGHT_GRAY},
    {"text": "  讲解语法和词汇 ✓", "size": 12, "color": GREEN},
]
add_multi_text(slide, Inches(7.2), Inches(2.1), Inches(5), Inches(4.5),
               qwen_cases, line_spacing=1.2)

add_slide_number(slide, 11)

# 第12页：通义千问优缺点+使用步骤
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_decorative_elements(slide)
add_section_header(slide, 6, "通义千问 — 优缺点与使用步骤")

add_rounded_rect(slide, Inches(0.8), Inches(1.2), Inches(3.5), Inches(5.8), BG_CARD, GREEN, Pt(2))
add_text_box(slide, Inches(1.0), Inches(1.4), Inches(3), Inches(0.5),
             "✅ 优点", font_size=22, color=GREEN, bold=True)
add_line(slide, Inches(1.0), Inches(1.9), Inches(3), GREEN, Pt(1))
qwen_pros = [
    "• 阿里生态，稳定可靠",
    "• 长文本处理能力强",
    "• 图像识别理解好",
    "• 与钉钉/阿里云打通",
    "• 免费额度充足",
    "• 支持文件上传分析",
    "• 中文理解优秀",
    "• 多模态能力强",
]
add_multi_text(slide, Inches(1.0), Inches(2.1), Inches(3), Inches(4.5),
               [{"text": p, "size": 14, "color": WHITE} for p in qwen_pros], line_spacing=1.5)

add_rounded_rect(slide, Inches(4.7), Inches(1.2), Inches(3.5), Inches(5.8), BG_CARD, RED, Pt(2))
add_text_box(slide, Inches(4.9), Inches(1.4), Inches(3), Inches(0.5),
             "⚠️ 缺点", font_size=22, color=RED, bold=True)
add_line(slide, Inches(4.9), Inches(1.9), Inches(3), RED, Pt(1))
qwen_cons = [
    "• 创意写作偏正式",
    "• 部分功能需登录",
    "• 实时信息有延迟",
    "• 数学计算偶有误",
    "• 界面偏商务风",
    "• 移动端功能有限",
    "• 高级模型需排队",
    "• 个性化不足",
]
add_multi_text(slide, Inches(4.9), Inches(2.1), Inches(3), Inches(4.5),
               [{"text": c, "size": 14, "color": WHITE} for c in qwen_cons], line_spacing=1.5)

add_rounded_rect(slide, Inches(8.6), Inches(1.2), Inches(3.9), Inches(5.8), BG_CARD, ORANGE, Pt(2))
add_text_box(slide, Inches(8.8), Inches(1.4), Inches(3.5), Inches(0.5),
             "📋 使用步骤", font_size=22, color=ORANGE, bold=True)
add_line(slide, Inches(8.8), Inches(1.9), Inches(3.5), ORANGE, Pt(1))
qwen_steps = [
    {"text": "Step 1：访问平台", "size": 15, "color": YELLOW, "bold": True},
    {"text": "  tongyi.aliyun.com", "size": 12, "color": LIGHT_GRAY},
    {"text": "", "size": 6},
    {"text": "Step 2：登录账号", "size": 15, "color": YELLOW, "bold": True},
    {"text": "  支付宝/淘宝/手机号", "size": 12, "color": LIGHT_GRAY},
    {"text": "", "size": 6},
    {"text": "Step 3：选择能力", "size": 15, "color": YELLOW, "bold": True},
    {"text": "  对话/文档/图片/代码", "size": 12, "color": LIGHT_GRAY},
    {"text": "", "size": 6},
    {"text": "Step 4：上传文件", "size": 15, "color": YELLOW, "bold": True},
    {"text": "  可上传PDF/Word/图片", "size": 12, "color": LIGHT_GRAY},
    {"text": "", "size": 6},
    {"text": "Step 5：提问分析", "size": 15, "color": YELLOW, "bold": True},
    {"text": "  输入具体问题获取", "size": 12, "color": LIGHT_GRAY},
    {"text": "  AI分析结果", "size": 12, "color": LIGHT_GRAY},
]
add_multi_text(slide, Inches(8.8), Inches(2.1), Inches(3.5), Inches(4.5),
               qwen_steps, line_spacing=1.2)

add_slide_number(slide, 12)

# ============================================================
# 第13-14页：剪映
# ============================================================
# 第13页：剪映介绍+案例
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_decorative_elements(slide)
add_section_header(slide, 7, "剪映 — AI视频创作工具")

add_rounded_rect(slide, Inches(0.8), Inches(1.2), Inches(5.8), Inches(5.8), BG_CARD, RED, Pt(2))
add_text_box(slide, Inches(1.0), Inches(1.4), Inches(5), Inches(0.5),
             "🎬 剪映", font_size=28, color=RED, bold=True)
add_line(slide, Inches(1.0), Inches(1.95), Inches(5), RED, Pt(1))

jy_info = [
    {"text": "开发商：字节跳动", "size": 15, "color": WHITE, "bold": True},
    {"text": "", "size": 8},
    {"text": "平台：电脑端 / 手机端", "size": 14, "color": LIGHT_GRAY},
    {"text": "", "size": 8},
    {"text": "核心功能：", "size": 15, "color": RED, "bold": True},
    {"text": "  • AI智能剪辑", "size": 13, "color": LIGHT_GRAY},
    {"text": "  • 自动字幕生成", "size": 13, "color": LIGHT_GRAY},
    {"text": "  • AI文案/脚本生成", "size": 13, "color": LIGHT_GRAY},
    {"text": "  • 智能美颜/滤镜", "size": 13, "color": LIGHT_GRAY},
    {"text": "  • 图文成片（文字转视频）", "size": 13, "color": LIGHT_GRAY},
    {"text": "  • 背景音乐智能匹配", "size": 13, "color": LIGHT_GRAY},
    {"text": "  • 模板一键成片", "size": 13, "color": LIGHT_GRAY},
]
add_multi_text(slide, Inches(1.0), Inches(2.1), Inches(5.3), Inches(4.5),
               jy_info, line_spacing=1.3)

add_rounded_rect(slide, Inches(7.0), Inches(1.2), Inches(5.5), Inches(5.8), BG_CARD, CYAN, Pt(2))
add_text_box(slide, Inches(7.2), Inches(1.4), Inches(5), Inches(0.5),
             "💡 实用案例", font_size=22, color=CYAN, bold=True)
add_line(slide, Inches(7.2), Inches(1.95), Inches(5), CYAN, Pt(1))

jy_cases = [
    {"text": "案例1：课堂汇报视频", "size": 15, "color": YELLOW, "bold": True},
    {"text": "作业：制作3分钟产品", "size": 12, "color": LIGHT_GRAY},
    {"text": "  介绍短视频", "size": 12, "color": LIGHT_GRAY},
    {"text": "剪映：导入素材→AI自动", "size": 12, "color": LIGHT_GRAY},
    {"text": "  剪辑→加字幕→导出 ✓", "size": 12, "color": GREEN},
    {"text": "", "size": 6},
    {"text": "案例2：图文成片", "size": 15, "color": YELLOW, "bold": True},
    {"text": "学生：写好产品介绍文案", "size": 12, "color": LIGHT_GRAY},
    {"text": "剪映：AI自动匹配素材+", "size": 12, "color": LIGHT_GRAY},
    {"text": "  配音+字幕=完整视频 ✓", "size": 12, "color": GREEN},
    {"text": "", "size": 6},
    {"text": "案例3：直播切片", "size": 15, "color": YELLOW, "bold": True},
    {"text": "电商实训：从1小时直播", "size": 12, "color": LIGHT_GRAY},
    {"text": "  中剪辑精彩片段", "size": 12, "color": LIGHT_GRAY},
    {"text": "剪映：AI识别高光时刻+", "size": 12, "color": LIGHT_GRAY},
    {"text": "  自动剪辑+包装 ✓", "size": 12, "color": GREEN},
]
add_multi_text(slide, Inches(7.2), Inches(2.1), Inches(5), Inches(4.5),
               jy_cases, line_spacing=1.2)

add_slide_number(slide, 13)

# 第14页：剪映优缺点+使用步骤
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_decorative_elements(slide)
add_section_header(slide, 7, "剪映 — 优缺点与使用步骤")

add_rounded_rect(slide, Inches(0.8), Inches(1.2), Inches(3.5), Inches(5.8), BG_CARD, GREEN, Pt(2))
add_text_box(slide, Inches(1.0), Inches(1.4), Inches(3), Inches(0.5),
             "✅ 优点", font_size=22, color=GREEN, bold=True)
add_line(slide, Inches(1.0), Inches(1.9), Inches(3), GREEN, Pt(1))
jy_pros = [
    "• 免费使用，功能强大",
    "• 操作简单，门槛低",
    "• AI自动字幕准确率高",
    "• 模板丰富，出片快",
    "• 与抖音无缝对接",
    "• 支持4K导出",
    "• 素材库丰富",
    "• 电脑/手机同步",
]
add_multi_text(slide, Inches(1.0), Inches(2.1), Inches(3), Inches(4.5),
               [{"text": p, "size": 14, "color": WHITE} for p in jy_pros], line_spacing=1.5)

add_rounded_rect(slide, Inches(4.7), Inches(1.2), Inches(3.5), Inches(5.8), BG_CARD, RED, Pt(2))
add_text_box(slide, Inches(4.9), Inches(1.4), Inches(3), Inches(0.5),
             "⚠️ 缺点", font_size=22, color=RED, bold=True)
add_line(slide, Inches(4.9), Inches(1.9), Inches(3), RED, Pt(1))
jy_cons = [
    "• 专业级特效有限",
    "• 高级功能需VIP",
    "• 大文件处理较慢",
    "• 自定义程度不如PR",
    "• 部分素材有版权",
    "• 导出有平台水印",
    "• 多轨道编辑不便",
    "• 色彩管理不专业",
]
add_multi_text(slide, Inches(4.9), Inches(2.1), Inches(3), Inches(4.5),
               [{"text": c, "size": 14, "color": WHITE} for c in jy_cons], line_spacing=1.5)

add_rounded_rect(slide, Inches(8.6), Inches(1.2), Inches(3.9), Inches(5.8), BG_CARD, RED, Pt(2))
add_text_box(slide, Inches(8.8), Inches(1.4), Inches(3.5), Inches(0.5),
             "📋 使用步骤", font_size=22, color=RED, bold=True)
add_line(slide, Inches(8.8), Inches(1.9), Inches(3.5), RED, Pt(1))
jy_steps = [
    {"text": "Step 1：下载安装", "size": 15, "color": YELLOW, "bold": True},
    {"text": "  电脑端/手机端均可", "size": 12, "color": LIGHT_GRAY},
    {"text": "", "size": 6},
    {"text": "Step 2：新建项目", "size": 15, "color": YELLOW, "bold": True},
    {"text": "  导入视频/图片素材", "size": 12, "color": LIGHT_GRAY},
    {"text": "", "size": 6},
    {"text": "Step 3：AI剪辑", "size": 15, "color": YELLOW, "bold": True},
    {"text": "  使用\"图文成片\"或", "size": 12, "color": LIGHT_GRAY},
    {"text": "  \"智能剪辑\"功能", "size": 12, "color": LIGHT_GRAY},
    {"text": "", "size": 6},
    {"text": "Step 4：添加元素", "size": 15, "color": YELLOW, "bold": True},
    {"text": "  字幕/音乐/特效/转场", "size": 12, "color": LIGHT_GRAY},
    {"text": "", "size": 6},
    {"text": "Step 5：导出分享", "size": 15, "color": YELLOW, "bold": True},
    {"text": "  选择分辨率导出，", "size": 12, "color": LIGHT_GRAY},
    {"text": "  可直接发抖音", "size": 12, "color": LIGHT_GRAY},
]
add_multi_text(slide, Inches(8.8), Inches(2.1), Inches(3.5), Inches(4.5),
               jy_steps, line_spacing=1.2)

add_slide_number(slide, 14)

# ============================================================
# 第15页：工具对比总览
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_decorative_elements(slide)
add_section_header(slide, 8, "五大AI工具对比总览")

# 表格头部
headers = ["功能/工具", "📱 豆包", "🧠 DeepSeek", "🎨 醒图", "💬 通义千问", "🎬 剪映"]
header_colors = [WHITE, CYAN, RGBColor(0x64, 0xFF, 0xDA), GREEN, ORANGE, RED]

# 表头
x = Inches(0.5)
for i, (h, hc) in enumerate(zip(headers, header_colors)):
    w = Inches(2.2) if i == 0 else Inches(2.0)
    add_rounded_rect(slide, x, Inches(1.2), w, Inches(0.5), BG_CARD, hc, Pt(1))
    add_text_box(slide, x, Inches(1.25), w, Inches(0.4),
                 h, font_size=13, color=hc, bold=True, alignment=PP_ALIGN.CENTER)
    x += w + Inches(0.05)

# 表格数据
rows = [
    ["核心定位", "AI对话", "深度推理", "AI修图", "多模态AI", "AI视频"],
    ["学习辅导", "★★★★", "★★★★★", "★", "★★★★", "★★"],
    ["创意设计", "★★★", "★★", "★★★★★", "★★★", "★★★★"],
    ["视频制作", "★", "★", "★★", "★★", "★★★★★"],
    ["文档处理", "★★★★", "★★★★★", "★", "★★★★★", "★★"],
    ["上手难度", "⭐低", "⭐⭐中", "⭐低", "⭐⭐中", "⭐低"],
    ["免费额度", "充足", "很充足", "基础免费", "充足", "基础免费"],
    ["最佳场景", "日常问答", "深度学习", "图片处理", "文件分析", "短视频"],
]

y = Inches(1.85)
for row in rows:
    x = Inches(0.5)
    for i, cell in enumerate(row):
        w = Inches(2.2) if i == 0 else Inches(2.0)
        color = WHITE if i == 0 else LIGHT_GRAY
        size = 13 if i == 0 else 12
        add_rounded_rect(slide, x, y, w, Inches(0.55), BG_CARD, RGBColor(0x1A, 0x22, 0x55), Pt(0.5))
        add_text_box(slide, x, y + Inches(0.05), w, Inches(0.45),
                     cell, font_size=size, color=color, alignment=PP_ALIGN.CENTER)
        x += w + Inches(0.05)
    y += Inches(0.6)

# 底部建议
add_rounded_rect(slide, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.7), BG_CARD, CYAN, Pt(1))
add_text_box(slide, Inches(0.7), Inches(6.55), Inches(12), Inches(0.6),
             "💡 选择建议：学习问答→豆包/通义千问 ｜ 深度学习→DeepSeek ｜ 图片处理→醒图 ｜ 视频制作→剪映",
             font_size=14, color=CYAN, alignment=PP_ALIGN.CENTER)

add_slide_number(slide, 15)

# ============================================================
# 第16页：提示词工程基础
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_decorative_elements(slide)
add_section_header(slide, 8, "AI工具核心技能：提示词工程")

# 左侧：什么是提示词
add_rounded_rect(slide, Inches(0.8), Inches(1.2), Inches(5.8), Inches(2.8), BG_CARD, CYAN, Pt(2))
add_text_box(slide, Inches(1.0), Inches(1.4), Inches(5.3), Inches(0.5),
             "🔑 什么是提示词（Prompt）？", font_size=20, color=CYAN, bold=True)
add_line(slide, Inches(1.0), Inches(1.9), Inches(5.3), CYAN, Pt(1))
prompt_def = [
    "提示词 = 你给AI的「指令」",
    "",
    "好的提示词 → 好的回答",
    "差的提示词 → 差的回答",
    "",
    "👉 提示词工程 = 学会如何向AI提问",
]
add_multi_text(slide, Inches(1.0), Inches(2.1), Inches(5.3), Inches(1.8),
               [{"text": p, "size": 15, "color": WHITE} for p in prompt_def], line_spacing=1.4)

# 右侧：好vs差对比
add_rounded_rect(slide, Inches(0.8), Inches(4.2), Inches(5.8), Inches(2.8), BG_CARD, GREEN, Pt(2))
add_text_box(slide, Inches(1.0), Inches(4.4), Inches(5.3), Inches(0.5),
             "📊 好提示词 vs 差提示词", font_size=18, color=GREEN, bold=True)
add_line(slide, Inches(1.0), Inches(4.9), Inches(5.3), GREEN, Pt(1))
compare = [
    {"text": "❌ 差：\"帮我写作文\"", "size": 13, "color": RED, "bold": True},
    {"text": "✅ 好：\"写一篇600字记叙文，", "size": 13, "color": GREEN, "bold": True},
    {"text": "  主题是\"难忘的一天\"，", "size": 13, "color": GREEN},
    {"text": "  用第一人称，要有细节描写\"", "size": 13, "color": GREEN},
    {"text": "", "size": 6},
    {"text": "❌ 差：\"解释光合作用\"", "size": 13, "color": RED, "bold": True},
    {"text": "✅ 好：\"用初中生能懂的", "size": 13, "color": GREEN, "bold": True},
    {"text": "  语言解释光合作用，", "size": 13, "color": GREEN},
    {"text": "  举一个生活中的例子\"", "size": 13, "color": GREEN},
]
add_multi_text(slide, Inches(1.0), Inches(5.0), Inches(5.3), Inches(1.8),
               compare, line_spacing=1.2)

# 底部：提示词公式
add_rounded_rect(slide, Inches(7.0), Inches(1.2), Inches(5.5), Inches(5.8), BG_CARD, MAGENTA, Pt(2))
add_text_box(slide, Inches(7.2), Inches(1.4), Inches(5), Inches(0.5),
             "📝 万能提示词公式", font_size=22, color=MAGENTA, bold=True)
add_line(slide, Inches(7.2), Inches(1.9), Inches(5), MAGENTA, Pt(1))

formula = [
    {"text": "角色 + 任务 + 要求 + 格式", "size": 20, "color": YELLOW, "bold": True, "align": PP_ALIGN.CENTER},
    {"text": "", "size": 10},
    {"text": "📌 角色：你是一位...", "size": 14, "color": WHITE},
    {"text": "   \"你是一位电商专业老师\"", "size": 12, "color": LIGHT_GRAY},
    {"text": "", "size": 8},
    {"text": "📌 任务：请帮我...", "size": 14, "color": WHITE},
    {"text": "   \"请帮我写一份产品推广文案\"", "size": 12, "color": LIGHT_GRAY},
    {"text": "", "size": 8},
    {"text": "📌 要求：字数/风格/重点...", "size": 14, "color": WHITE},
    {"text": "   \"200字，活泼风格，突出性价比\"", "size": 12, "color": LIGHT_GRAY},
    {"text": "", "size": 8},
    {"text": "📌 格式：表格/列表/段落...", "size": 14, "color": WHITE},
    {"text": "   \"用三点式列出核心卖点\"", "size": 12, "color": LIGHT_GRAY},
]
add_multi_text(slide, Inches(7.2), Inches(2.1), Inches(5), Inches(4.5),
               formula, line_spacing=1.3)

add_slide_number(slide, 16)

# ============================================================
# 第17页：提示词实战演练
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_decorative_elements(slide)
add_section_header(slide, 8, "提示词实战演练")

examples = [
    ("📚 学习场景", "你是中职电商专业老师，请帮我列出\"网店运营\"课程的5个核心知识点，每个知识点用1-2句话解释，最后给出学习建议。", CYAN),
    ("✍️ 写作场景", "你是一位文案策划，请为一款\"无线蓝牙耳机\"写3条小红书风格的推广文案，每条不超过100字，要活泼有趣，带emoji表情。", GREEN),
    ("🎯 规划场景", "我是一名中职高一学生，专业是跨境电商。请帮我制定一个学期学习计划，包含专业技能、文化课、课外活动三个方面，用表格形式呈现。", MAGENTA),
    ("🛒 电商场景", "你是一位电商运营专家，请帮我分析\"如何在抖音上推广一款国产护肤品\"，列出5个具体可行的推广策略，每个策略说明预期效果和投入成本。", ORANGE),
]

y = Inches(1.2)
for i, (title, prompt, color) in enumerate(examples):
    add_rounded_rect(slide, Inches(0.8), y, Inches(11.7), Inches(1.3), BG_CARD, color, Pt(1))
    add_text_box(slide, Inches(1.0), y + Inches(0.1), Inches(3), Inches(0.4),
                 title, font_size=16, color=color, bold=True)
    add_text_box(slide, Inches(1.0), y + Inches(0.5), Inches(11.3), Inches(0.7),
                 f"💬 \"{prompt}\"", font_size=12, color=LIGHT_GRAY)
    y += Inches(1.45)

add_slide_number(slide, 17)

# ============================================================
# 第18页：AI工具使用注意事项
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_decorative_elements(slide)
add_section_header(slide, 9, "AI工具使用注意事项 ⚠️")

# 四宫格
cards = [
    ("🚫 不要完全依赖AI", [
        "• AI只是辅助工具",
        "• 独立思考最重要",
        "• 不要直接抄AI答案",
        "• 理解比结果更重要",
        "• 培养自己的判断力",
    ], RED, Inches(0.8)),
    ("✅ 要验证AI答案", [
        "• 重要信息交叉验证",
        "• 查资料确认准确性",
        "• AI可能\"幻觉\"编造",
        "• 数学题要自己验算",
        "• 事实性内容查权威源",
    ], YELLOW, Inches(4.3)),
    ("🔒 注意隐私安全", [
        "• 不输入个人敏感信息",
        "• 不上传身份证等证件",
        "• 不泄露家庭住址电话",
        "• 注意账号密码安全",
        "• 了解平台隐私政策",
    ], CYAN, Inches(7.8)),
    ("⚖️ 遵守学术诚信", [
        "• 作业注明AI辅助部分",
        "• 考试不使用AI工具",
        "• 尊重知识产权",
        "• AI生成内容需审核",
        "• 培养诚信品质",
    ], GREEN, Inches(11.3)),
]

for title, items, color, x_pos in cards:
    add_rounded_rect(slide, x_pos, Inches(1.2), Inches(3.2), Inches(5.8), BG_CARD, color, Pt(2))
    add_text_box(slide, x_pos + Inches(0.2), Inches(1.4), Inches(2.8), Inches(0.5),
                 title, font_size=16, color=color, bold=True)
    add_line(slide, x_pos + Inches(0.2), Inches(1.95), Inches(2.8), color, Pt(1))
    add_multi_text(slide, x_pos + Inches(0.2), Inches(2.2), Inches(2.8), Inches(4.5),
                   [{"text": item, "size": 14, "color": WHITE} for item in items], line_spacing=1.6)

add_slide_number(slide, 18)

# ============================================================
# 第19页：课堂实践任务
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_decorative_elements(slide)
add_section_header(slide, 10, "课堂实践任务 🎯")

tasks = [
    ("任务一：AI对话体验（15分钟）", CYAN, [
        "📱 打开豆包或通义千问",
        "1. 用\"万能公式\"写一个提示词",
        "2. 向AI提问一个你感兴趣的问题",
        "3. 对比\"好提示词\"和\"差提示词\"的结果差异",
        "4. 记录：AI回答的质量如何？有什么不足？",
        "📝 提交：截图+200字体验报告",
    ]),
    ("任务二：AI修图实战（15分钟）", GREEN, [
        "🎨 打开醒图APP",
        "1. 选择一张自己的照片",
        "2. 使用AI美颜功能调整",
        "3. 使用\"一键抠图\"换背景",
        "4. 添加文字制作成个人海报",
        "📝 提交：原图vs成品图对比",
    ]),
    ("任务三：AI视频创作（15分钟）", RED, [
        "🎬 打开剪映",
        "1. 使用\"图文成片\"功能",
        "2. 输入一段产品介绍文案",
        "3. AI自动生成视频后手动调整",
        "4. 添加字幕和背景音乐",
        "📝 提交：导出视频文件",
    ]),
]

y = Inches(1.2)
for title, color, items in tasks:
    add_rounded_rect(slide, Inches(0.8), y, Inches(11.7), Inches(1.7), BG_CARD, color, Pt(2))
    add_text_box(slide, Inches(1.0), y + Inches(0.1), Inches(11), Inches(0.4),
                 title, font_size=18, color=color, bold=True)
    add_line(slide, Inches(1.0), y + Inches(0.5), Inches(11), color, Pt(1))
    add_multi_text(slide, Inches(1.0), y + Inches(0.6), Inches(11), Inches(1.0),
                   [{"text": item, "size": 14, "color": WHITE} for item in items], line_spacing=1.4)
    y += Inches(1.85)

add_slide_number(slide, 19)

# ============================================================
# 第20页：AI工具学习路线图
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_decorative_elements(slide)
add_section_header(slide, 10, "AI工具学习路线图 🗺️")

phases = [
    ("第一阶段\n入门", "第1-2周", [
        "• 下载注册5个AI工具",
        "• 体验基本对话功能",
        "• 理解什么是AI",
        "• 学会基础提示词",
    ], CYAN, Inches(0.8)),
    ("第二阶段\n进阶", "第3-4周", [
        "• 掌握提示词技巧",
        "• 学习AI修图实操",
        "• 尝试AI视频制作",
        "• 完成课堂实践任务",
    ], GREEN, Inches(3.5)),
    ("第三阶段\n应用", "第5-6周", [
        "• AI辅助专业课学习",
        "• 用AI完成课程作业",
        "• 尝试AI创意项目",
        "• 小组协作AI项目",
    ], MAGENTA, Inches(6.2)),
    ("第四阶段\n精通", "第7-8周", [
        "• 综合运用多工具",
        "• 独立完成AI作品",
        "• 参加AI技能比赛",
        "• 分享AI使用经验",
    ], ORANGE, Inches(8.9)),
]

for title, time, items, color, x_pos in phases:
    # 阶段卡片
    add_rounded_rect(slide, x_pos, Inches(1.2), Inches(2.3), Inches(5.5), BG_CARD, color, Pt(2))
    # 阶段标题
    add_text_box(slide, x_pos + Inches(0.1), Inches(1.4), Inches(2.1), Inches(0.7),
                 title, font_size=18, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    # 时间
    add_text_box(slide, x_pos + Inches(0.1), Inches(2.1), Inches(2.1), Inches(0.3),
                 time, font_size=12, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
    add_line(slide, x_pos + Inches(0.2), Inches(2.45), Inches(1.9), color, Pt(1))
    # 内容
    add_multi_text(slide, x_pos + Inches(0.15), Inches(2.7), Inches(2.0), Inches(3.5),
                   [{"text": item, "size": 13, "color": WHITE} for item in items], line_spacing=1.6)
    # 连接箭头
    if x_pos < Inches(8.9):
        add_text_box(slide, x_pos + Inches(2.3), Inches(3.5), Inches(0.5), Inches(0.5),
                     "→", font_size=24, color=color, bold=True, alignment=PP_ALIGN.CENTER)

# 底部总结
add_rounded_rect(slide, Inches(0.8), Inches(6.9), Inches(11.7), Inches(0.4), BG_CARD, CYAN, Pt(1))
add_text_box(slide, Inches(1.0), Inches(6.92), Inches(11.3), Inches(0.35),
             "💡 学习AI工具的关键：多练、多问、多对比、多总结",
             font_size=14, color=CYAN, alignment=PP_ALIGN.CENTER, bold=True)

add_slide_number(slide, 20)

# ============================================================
# 第21页：常见问题FAQ
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_decorative_elements(slide)
add_section_header(slide, 10, "常见问题 FAQ")

faqs = [
    ("Q1：AI工具都要付费吗？", "大部分AI工具有免费版本，足够学生日常使用。高级功能可能需要付费，但学生通常不需要。"),
    ("Q2：用AI写作业算作弊吗？", "取决于学校规定。建议：将AI作为学习辅助（如解释概念、提供思路），而非直接抄答案。重要作业应注明AI辅助部分。"),
    ("Q3：AI回答一定准确吗？", "不一定！AI可能产生\"幻觉\"（编造信息）。重要信息一定要交叉验证，不要盲目相信。"),
    ("Q4：哪个AI工具最适合学生？", "入门推荐豆包（中文好、易上手）；深度学习推荐DeepSeek（推理强）；修图用醒图；视频用剪映。"),
    ("Q5：AI会取代我的工作吗？", "AI会取代\"不会用AI的人\"。学会使用AI工具，反而能提升你的竞争力。关键是成为\"会用AI的人\"。"),
]

y = Inches(1.2)
for q, a in faqs:
    add_rounded_rect(slide, Inches(0.8), y, Inches(11.7), Inches(1.0), BG_CARD, CYAN_DIM, Pt(1))
    add_text_box(slide, Inches(1.0), y + Inches(0.05), Inches(11), Inches(0.35),
                 q, font_size=15, color=CYAN, bold=True)
    add_text_box(slide, Inches(1.0), y + Inches(0.4), Inches(11), Inches(0.55),
                 a, font_size=13, color=LIGHT_GRAY)
    y += Inches(1.1)

add_slide_number(slide, 21)

# ============================================================
# 第22页：结束页
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

# 背景装饰
add_circle(slide, Inches(-2), Inches(-2), Inches(6), RGBColor(0x0D, 0x14, 0x3A))
add_circle(slide, Inches(10), Inches(4), Inches(5), RGBColor(0x0D, 0x14, 0x3A))

add_line(slide, Inches(1.5), Inches(2.5), Inches(10.3), CYAN, Pt(2))
add_line(slide, Inches(1.5), Inches(5.5), Inches(10.3), MAGENTA, Pt(1))

add_text_box(slide, Inches(1.5), Inches(2.8), Inches(10.3), Inches(1.0),
             "谢谢观看！", font_size=54, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1.5), Inches(4.0), Inches(10.3), Inches(0.6),
             "AI时代，学会与AI共舞 🕺", font_size=26, color=CYAN, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1.5), Inches(4.8), Inches(10.3), Inches(0.5),
             "掌握AI工具 = 掌握未来竞争力", font_size=20, color=MAGENTA, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1.5), Inches(5.7), Inches(10.3), Inches(0.4),
             "中职高一信息技术课 · AI工具使用指南", font_size=16, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

add_circle(slide, Inches(1.8), Inches(5.8), Inches(0.1), CYAN)
add_circle(slide, Inches(11.4), Inches(5.8), Inches(0.1), MAGENTA)

add_slide_number(slide, 22)

# ===== 保存 =====
output_path = "/home/admin/.openclaw/workspace/papers/AI工具使用指南_中职高一.pptx"
prs.save(output_path)
print(f"PPT已保存: {output_path}")
print(f"总页数: {len(prs.slides)}")
