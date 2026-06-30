#!/usr/bin/env python3
"""
生成「智能体赋能高校教学新范式」直播课程汇报 PPTX
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── 配色方案 ──
DARK_BG    = RGBColor(0x0B, 0x11, 0x20)
BLUE_ACCENT = RGBColor(0x3B, 0x82, 0xF6)
LIGHT_BLUE  = RGBColor(0x60, 0xA5, 0xFA)
WHITE       = RGBColor(0xF0, 0xF4, 0xF8)
MUTED       = RGBColor(0x94, 0xA3, 0xB8)
CARD_BG     = RGBColor(0x1E, 0x29, 0x3B)

def add_bg(slide, color=DARK_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, fill_color, border_color=None, border_width=None):
    from pptx.util import Emu
    from pptx.enum.shapes import MSO_SHAPE
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(border_width or 1)
    else:
        shape.line.fill.background()
    # Round corners
    shape.adjustments[0] = 0.1
    return shape

def set_text(tf, text, size=24, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="微软雅黑"):
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment

def add_multiline(tf, lines, size=24, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, line_spacing=1.4, font_name="微软雅黑"):
    """Add multiple paragraphs to a text frame."""
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font_name
        p.alignment = alignment
        p.space_after = Pt(size * (line_spacing - 1) * 4)

# ═══════════════════════════════════════════
# Slide 1: Cover
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide)

# Accent line at top
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), LIGHT_BLUE)

# Small label
txBox = slide.shapes.add_textbox(Inches(1.5), Inches(1.2), Inches(10), Inches(0.6))
set_text(txBox.text_frame, "直 播 课 程", 22, LIGHT_BLUE, bold=True, alignment=PP_ALIGN.LEFT)

# Main title
txBox = slide.shapes.add_textbox(Inches(1.5), Inches(2.0), Inches(10), Inches(1.8))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "智能体赋能高校教学新范式"
p.font.size = Pt(60)
p.font.color.rgb = WHITE
p.font.bold = True
p.font.name = "微软雅黑"
p.alignment = PP_ALIGN.CENTER

# Subtitle
txBox = slide.shapes.add_textbox(Inches(1.5), Inches(3.8), Inches(10), Inches(0.8))
set_text(txBox.text_frame, "小龙虾 + Manus 一站式解决方案", 36, BLUE_ACCENT, bold=True, alignment=PP_ALIGN.CENTER)

# Divider line
txBox = slide.shapes.add_textbox(Inches(6.2), Inches(4.8), Inches(1), Inches(0.05))
add_shape(slide, Inches(6.2), Inches(4.8), Inches(1), Inches(0.03), BLUE_ACCENT)

# Info row
txBox = slide.shapes.add_textbox(Inches(1.5), Inches(5.2), Inches(10), Inches(1.0))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "📅 6月27日（周五） 15:00-16:00"
p.font.size = Pt(30)
p.font.color.rgb = MUTED
p.font.name = "微软雅黑"
p.alignment = PP_ALIGN.CENTER

# Book reference
txBox = slide.shapes.add_textbox(Inches(1.5), Inches(6.2), Inches(10), Inches(0.6))
set_text(txBox.text_frame, "📘 《Manus智能体全攻略》· 清华大学出版社", 24, MUTED, alignment=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════
# Slide 2: Book Introduction
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), LIGHT_BLUE)

txBox = slide.shapes.add_textbox(Inches(1.2), Inches(0.5), Inches(11), Inches(0.5))
set_text(txBox.text_frame, "所 用 图 书", 20, LIGHT_BLUE, bold=True)

txBox = slide.shapes.add_textbox(Inches(1.2), Inches(1.2), Inches(11), Inches(1.2))
set_text(txBox.text_frame, "《Manus智能体全攻略》", 52, WHITE, bold=True)

# Info cards on left
info_items = [
    ("出 版 社", "清华大学出版社"),
    ("作  者", "诸葛斌 等"),
    ("定  位", "国内首本智能体教学实战指南"),
]
for i, (label, value) in enumerate(info_items):
    y = Inches(2.6) + Inches(i * 1.3)
    txBox = slide.shapes.add_textbox(Inches(1.2), y, Inches(5), Inches(0.4))
    set_text(txBox.text_frame, label, 20, MUTED)
    txBox = slide.shapes.add_textbox(Inches(1.2), y + Inches(0.4), Inches(5), Inches(0.6))
    set_text(txBox.text_frame, value, 30, WHITE, bold=True)

# Book cover mockup (right side)
card = add_shape(slide, Inches(8.5), Inches(2.2), Inches(3.2), Inches(4.2), RGBColor(0x1E, 0x3A, 0x5F), BLUE_ACCENT, 2)
txBox = slide.shapes.add_textbox(Inches(8.7), Inches(3.0), Inches(2.8), Inches(2.5))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "📘"
p.font.size = Pt(72)
p.alignment = PP_ALIGN.CENTER
p = tf.add_paragraph()
p.text = "Manus智能体\n全攻略"
p.font.size = Pt(28)
p.font.color.rgb = WHITE
p.font.bold = True
p.font.name = "微软雅黑"
p.alignment = PP_ALIGN.CENTER
p.space_before = Pt(16)
p = tf.add_paragraph()
p.text = "清华大学出版社"
p.font.size = Pt(16)
p.font.color.rgb = MUTED
p.font.name = "微软雅黑"
p.alignment = PP_ALIGN.CENTER

# ═══════════════════════════════════════════
# Slide 3: Broadcast Content
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), LIGHT_BLUE)

txBox = slide.shapes.add_textbox(Inches(1.2), Inches(0.5), Inches(11), Inches(0.5))
set_text(txBox.text_frame, "直 播 内 容", 20, LIGHT_BLUE, bold=True)

txBox = slide.shapes.add_textbox(Inches(1.2), Inches(1.1), Inches(11), Inches(1.0))
set_text(txBox.text_frame, "解读 + 实操，一次搞定", 44, WHITE, bold=True)

txBox = slide.shapes.add_textbox(Inches(1.2), Inches(2.0), Inches(11), Inches(0.8))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "解读清华版《Manus智能体全攻略》，现场演示五大教学场景，"
p.font.size = Pt(24)
p.font.color.rgb = MUTED
p.font.name = "微软雅黑"
p = tf.add_paragraph()
p.text = "实操 2 小时教案制作压缩至 10 分钟"
p.font.size = Pt(28)
p.font.color.rgb = LIGHT_BLUE
p.font.bold = True
p.font.name = "微软雅黑"

# 4 content cards
cards_data = [
    ("📚", "课件智能生成", '"小龙虾三部曲"精品课件'),
    ("🔬", "教学案例开发", "烟草数据挖掘、网络课程动画等"),
    ("📝", "论文协作写作", "从选题到 IEEE 成稿全流程"),
    ("🎬", "教学视频制作", "PPT 自动转教学视频"),
]
for i, (icon, title, desc) in enumerate(cards_data):
    col = i % 2
    row = i // 2
    x = Inches(1.2) + col * Inches(5.8)
    y = Inches(3.3) + row * Inches(2.0)
    card = add_shape(slide, x, y, Inches(5.5), Inches(1.7), CARD_BG, RGBColor(0x33, 0x41, 0x55), 1)
    
    # Icon
    txBox = slide.shapes.add_textbox(x + Inches(0.3), y + Inches(0.2), Inches(0.8), Inches(0.8))
    set_text(txBox.text_frame, icon, 40, WHITE)
    
    # Title
    txBox = slide.shapes.add_textbox(x + Inches(1.2), y + Inches(0.2), Inches(4), Inches(0.6))
    set_text(txBox.text_frame, title, 26, WHITE, bold=True)
    
    # Desc
    txBox = slide.shapes.add_textbox(x + Inches(1.2), y + Inches(0.85), Inches(4), Inches(0.7))
    set_text(txBox.text_frame, desc, 20, MUTED)

# ═══════════════════════════════════════════
# Slide 4: 5 Scenarios
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), LIGHT_BLUE)

txBox = slide.shapes.add_textbox(Inches(1.2), Inches(0.5), Inches(11), Inches(0.5))
set_text(txBox.text_frame, "五 大 教 学 场 景", 20, LIGHT_BLUE, bold=True)

txBox = slide.shapes.add_textbox(Inches(1.2), Inches(1.1), Inches(11), Inches(0.8))
set_text(txBox.text_frame, "从课件到视频，全流程覆盖", 44, WHITE, bold=True)

scenarios = [
    ("01", "课件智能生成", '"小龙虾三部曲"精品课件自动化生成'),
    ("02", "教学案例开发", "烟草数据挖掘案例 · 网络课程动画 · 微信小程序实战"),
    ("03", "论文协作写作", "从选题到 IEEE 成稿，全流程 AI 辅助"),
    ("04", "教学视频制作", "PPT 自动转教学视频，无需专业剪辑"),
    ("05", "教案快速制作", "2 小时教案制作压缩至 10 分钟"),
]

for i, (num, title, desc) in enumerate(scenarios):
    y = Inches(2.2) + Inches(i * 1.05)
    card = add_shape(slide, Inches(1.2), y, Inches(11), Inches(0.9), CARD_BG, RGBColor(0x33, 0x41, 0x55), 1)
    
    # Number
    txBox = slide.shapes.add_textbox(Inches(1.5), y + Inches(0.15), Inches(0.8), Inches(0.6))
    set_text(txBox.text_frame, num, 24, BLUE_ACCENT, bold=True)
    
    # Title
    txBox = slide.shapes.add_textbox(Inches(2.4), y + Inches(0.1), Inches(4), Inches(0.5))
    set_text(txBox.text_frame, title, 28, WHITE, bold=True)
    
    # Desc
    txBox = slide.shapes.add_textbox(Inches(2.4), y + Inches(0.52), Inches(9), Inches(0.4))
    set_text(txBox.text_frame, desc, 20, MUTED)

# ═══════════════════════════════════════════
# Slide 5: Author
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), LIGHT_BLUE)

txBox = slide.shapes.add_textbox(Inches(1.2), Inches(0.5), Inches(11), Inches(0.5))
set_text(txBox.text_frame, "主 讲 人", 20, LIGHT_BLUE, bold=True)

# Avatar mockup
card = add_shape(slide, Inches(1.2), Inches(1.5), Inches(2.6), Inches(3.5), RGBColor(0x1E, 0x3A, 0x5F), BLUE_ACCENT, 2)
txBox = slide.shapes.add_textbox(Inches(1.2), Inches(2.5), Inches(2.6), Inches(1.0))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "👤"
p.font.size = Pt(72)
p.alignment = PP_ALIGN.CENTER
p = tf.add_paragraph()
p.text = "诸葛斌"
p.font.size = Pt(28)
p.font.color.rgb = WHITE
p.font.bold = True
p.font.name = "微软雅黑"
p.alignment = PP_ALIGN.CENTER
p.space_before = Pt(8)
p = tf.add_paragraph()
p.text = "教授"
p.font.size = Pt(20)
p.font.color.rgb = MUTED
p.font.name = "微软雅黑"
p.alignment = PP_ALIGN.CENTER

# Author info
txBox = slide.shapes.add_textbox(Inches(4.5), Inches(1.5), Inches(8), Inches(0.8))
set_text(txBox.text_frame, "诸葛斌 教授", 48, WHITE, bold=True)

txBox = slide.shapes.add_textbox(Inches(4.5), Inches(2.4), Inches(8), Inches(0.6))
set_text(txBox.text_frame, "浙江工商大学 · 信息与电子工程学院 / 萨塞克斯人工智能学院", 24, BLUE_ACCENT)

achievements = [
    "研究方向：互联网应用开发与 AI 教育",
    "获浙江省技术发明一等奖",
    "联合阿里钉钉撰写国内首本低代码开发教材",
    "获 2025 全国高校人工智能教育大会优秀案例一等奖",
]
for i, item in enumerate(achievements):
    y = Inches(3.3) + Inches(i * 0.7)
    # Blue bar
    add_shape(slide, Inches(4.5), y, Inches(0.08), Inches(0.5), BLUE_ACCENT)
    txBox = slide.shapes.add_textbox(Inches(4.8), y - Inches(0.05), Inches(7), Inches(0.6))
    set_text(txBox.text_frame, item, 22, MUTED)

# ═══════════════════════════════════════════
# Slide 6: Benefits
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), LIGHT_BLUE)

txBox = slide.shapes.add_textbox(Inches(1.2), Inches(0.5), Inches(11), Inches(0.5))
set_text(txBox.text_frame, "入 群 福 利", 20, LIGHT_BLUE, bold=True)

txBox = slide.shapes.add_textbox(Inches(1.2), Inches(1.1), Inches(11), Inches(0.8))
set_text(txBox.text_frame, "扫码加入读者服务群", 44, WHITE, bold=True)

txBox = slide.shapes.add_textbox(Inches(1.2), Inches(1.9), Inches(11), Inches(0.5))
set_text(txBox.text_frame, "此群长期有效，群成员可获得以下福利", 22, MUTED)

benefits = [
    ("🤖", "小龙虾 AI 体验", "群内部署 OpenClaw 智能体，实时体验课件生成、教案制作等 AI 能力"),
    ("📦", "教学资料包", "小龙虾三部曲课件、数据挖掘案例、16 章教学动画、微信小程序案例等全套 Manus 实战案例"),
    ("🎁", "免费样书抽奖", "直播间专享 10 本《Manus智能体全攻略》免费样书赠送（名额有限，抽奖获得）"),
    ("💬", "教学交流社区", '"智能体"系列课程教学交流社区——高校教师 AI 教学实践交流、问题解答、经验分享'),
]
for i, (icon, title, desc) in enumerate(benefits):
    col = i % 2
    row = i // 2
    x = Inches(1.2) + col * Inches(5.8)
    y = Inches(2.6) + row * Inches(2.3)
    card = add_shape(slide, x, y, Inches(5.5), Inches(2.0), CARD_BG, RGBColor(0x33, 0x41, 0x55), 1)
    
    txBox = slide.shapes.add_textbox(x + Inches(0.3), y + Inches(0.2), Inches(0.8), Inches(0.8))
    set_text(txBox.text_frame, icon, 40, WHITE)
    
    txBox = slide.shapes.add_textbox(x + Inches(1.2), y + Inches(0.2), Inches(4), Inches(0.5))
    set_text(txBox.text_frame, title, 26, WHITE, bold=True)
    
    txBox = slide.shapes.add_textbox(x + Inches(1.2), y + Inches(0.8), Inches(4), Inches(1.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = desc
    p.font.size = Pt(18)
    p.font.color.rgb = MUTED
    p.font.name = "微软雅黑"

# ═══════════════════════════════════════════
# Slide 7: Closing
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

# Accent border
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), LIGHT_BLUE)
add_shape(slide, Inches(0), Inches(7.44), Inches(13.333), Inches(0.06), LIGHT_BLUE)

# Label
txBox = slide.shapes.add_textbox(Inches(1.5), Inches(1.2), Inches(10), Inches(0.5))
set_text(txBox.text_frame, "敬 请 期 待", 22, LIGHT_BLUE, bold=True, alignment=PP_ALIGN.CENTER)

# Title
txBox = slide.shapes.add_textbox(Inches(1.5), Inches(2.0), Inches(10), Inches(1.4))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "智能体赋能高校教学新范式"
p.font.size = Pt(56)
p.font.color.rgb = WHITE
p.font.bold = True
p.font.name = "微软雅黑"
p.alignment = PP_ALIGN.CENTER

# Subtitle
txBox = slide.shapes.add_textbox(Inches(1.5), Inches(3.4), Inches(10), Inches(0.7))
set_text(txBox.text_frame, "小龙虾 + Manus 一站式解决方案", 36, BLUE_ACCENT, bold=True, alignment=PP_ALIGN.CENTER)

# Info row
txBox = slide.shapes.add_textbox(Inches(1.5), Inches(4.4), Inches(10), Inches(0.6))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "📅 6月27日（周五）  ·  ⏰ 15:00-16:00  ·  👤 诸葛斌 教授"
p.font.size = Pt(28)
p.font.color.rgb = MUTED
p.font.name = "微软雅黑"
p.alignment = PP_ALIGN.CENTER

# CTA box
card = add_shape(slide, Inches(3.5), Inches(5.4), Inches(6.5), Inches(1.2), RGBColor(0x0E, 0x1E, 0x3F), BLUE_ACCENT, 2)
txBox = slide.shapes.add_textbox(Inches(3.8), Inches(5.55), Inches(5.8), Inches(0.5))
set_text(txBox.text_frame, "📱 请提前扫码加入读者服务群", 26, LIGHT_BLUE, bold=True, alignment=PP_ALIGN.CENTER)
txBox = slide.shapes.add_textbox(Inches(3.8), Inches(6.1), Inches(5.8), Inches(0.4))
set_text(txBox.text_frame, "享 4 大专属福利 · 群长期有效", 20, MUTED, alignment=PP_ALIGN.CENTER)

# ── 保存 ──
output_path = "/home/admin/.openclaw/workspace/lobster-network/直播课程汇报.pptx"
prs.save(output_path)
print(f"✅ PPTX saved: {output_path}")
print(f"   Slides: {len(prs.slides)}")
