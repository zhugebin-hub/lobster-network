#!/usr/bin/env python3
"""
直播课程汇报 PPTX - 布局优化版
修复文字遮挡、间距过小问题
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

DARK     = RGBColor(0x0B, 0x11, 0x20)
BLUE     = RGBColor(0x3B, 0x82, 0xF6)
LIGHT    = RGBColor(0x60, 0xA5, 0xFA)
WHITE    = RGBColor(0xF0, 0xF4, 0xF8)
MUTED    = RGBColor(0x94, 0xA3, 0xB8)
CARD     = RGBColor(0x1E, 0x29, 0x3B)
GREEN    = RGBColor(0x34, 0xD3, 0x99)
RED      = RGBColor(0xFF, 0x6B, 0x6B)
BORDER   = RGBColor(0x33, 0x41, 0x55)

def add_bg(slide, color=DARK):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = color

def card(s, l, t, w, h, fill=CARD, border=None, bw=None, r=0.1):
    shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if border:
        shape.line.color.rgb = border
        shape.line.width = Pt(bw or 1)
    else:
        shape.line.fill.background()
    shape.adjustments[0] = r
    return shape

def txt(s, l, t, w, h, text, size=18, color=WHITE, align=PP_ALIGN.LEFT, bold=False, sp=1.3):
    tx = s.shapes.add_textbox(l, t, w, h)
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "微软雅黑"
    p.alignment = align
    p.line_spacing = sp
    return tf

def bullets(s, l, t, w, h, items):
    tx = s.shapes.add_textbox(l, t, w, h)
    tf = tx.text_frame
    tf.word_wrap = True
    for i, (b, text, sz, clr) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = b + " " + text
        p.font.size = Pt(sz)
        p.font.color.rgb = clr
        p.font.name = "微软雅黑"
        p.space_after = Pt(6)
        p.line_spacing = 1.2
    return tf

# ============ SLIDE 1: COVER ============
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
card(s, Inches(0), Inches(0), Inches(13.333), Inches(0.05), LIGHT, r=0)
txt(s, Inches(1.2), Inches(0.8), Inches(5), Inches(0.4), "直播课程 - 清华大学出版社", 16, LIGHT, True)
txt(s, Inches(0.5), Inches(2.0), Inches(12.3), Inches(1.6), "智能体赋能高校教学新范式", 48, WHITE, align=PP_ALIGN.CENTER, bold=True, sp=1.15)
txt(s, Inches(1.0), Inches(3.5), Inches(11.3), Inches(0.6), "小龙虾 + Manus 一站式解决方案", 28, BLUE, align=PP_ALIGN.CENTER, bold=True)
card(s, Inches(1.2), Inches(4.4), Inches(11), Inches(0.9), CARD, BORDER, 1)
txt(s, Inches(1.5), Inches(4.5), Inches(3.3), Inches(0.7), "诸葛斌 教授", 20, WHITE, align=PP_ALIGN.CENTER, bold=True)
txt(s, Inches(5.0), Inches(4.5), Inches(3.3), Inches(0.7), "6月27日 周五", 20, WHITE, align=PP_ALIGN.CENTER, bold=True)
txt(s, Inches(8.5), Inches(4.5), Inches(3.3), Inches(0.7), "15:00 - 16:00", 20, WHITE, align=PP_ALIGN.CENTER, bold=True)
card(s, Inches(3.5), Inches(5.8), Inches(6.3), Inches(0.55), RGBColor(0x1E,0x3A,0x5F), BLUE, 2)
txt(s, Inches(3.7), Inches(5.82), Inches(5.9), Inches(0.5), "基于 Manus 智能体全攻略", 16, WHITE, align=PP_ALIGN.CENTER, bold=False)

# ============ SLIDE 2: PAINS vs SOLUTIONS ============
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
card(s, Inches(0), Inches(0), Inches(13.333), Inches(0.05), LIGHT, r=0)
txt(s, Inches(1.2), Inches(0.4), Inches(11), Inches(0.4), "背景与痛点", 16, LIGHT, True)
txt(s, Inches(1.2), Inches(0.9), Inches(11), Inches(0.6), "高校教师的 AI 教学困境", 32, WHITE, True)

pains = [("课件制作耗时","一份精品课件需数小时甚至数天"),("案例更新滞后","教学案例跟不上技术发展"),("论文写作困难","从选题到成稿周期长"),("视频制作复杂","需要专业剪辑技能")]
for i,(title,desc) in enumerate(pains):
    y = Inches(1.7) + Inches(i*1.1)
    card(s, Inches(0.6), y, Inches(5.6), Inches(0.95), CARD, RED, 1)
    txt(s, Inches(0.9), y+Inches(0.08), Inches(5), Inches(0.35), title, 18, RED, True)
    txt(s, Inches(0.9), y+Inches(0.45), Inches(5), Inches(0.35), desc, 14, MUTED)

txt(s, Inches(6.4), Inches(3.2), Inches(0.5), Inches(0.6), ">", 32, BLUE, align=PP_ALIGN.CENTER, bold=False)

sols = [("课件生成","2h教案 -> 10min"),("案例开发","自动挖掘+可视化"),("论文协作","选题到成稿全流程"),("视频制作","PPT自动转视频")]
for i,(title,desc) in enumerate(sols):
    y = Inches(1.7) + Inches(i*1.1)
    card(s, Inches(7.0), y, Inches(5.6), Inches(0.95), RGBColor(0x0E,0x3A,0x2F), GREEN, 1)
    txt(s, Inches(7.3), y+Inches(0.08), Inches(5), Inches(0.35), title, 18, GREEN, True)
    txt(s, Inches(7.3), y+Inches(0.45), Inches(5), Inches(0.35), desc, 14, MUTED)

# ============ SLIDE 3: SCENARIOS (UP) ============
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
card(s, Inches(0), Inches(0), Inches(13.333), Inches(0.05), LIGHT, r=0)
txt(s, Inches(1.2), Inches(0.35), Inches(11), Inches(0.4), "五大教学场景（上）", 16, LIGHT, True)

# Scene 1
card(s, Inches(0.5), Inches(0.9), Inches(12.3), Inches(1.7), CARD, BORDER, 1)
txt(s, Inches(0.8), Inches(1.0), Inches(0.7), Inches(0.5), "01", 28, LIGHT, align=PP_ALIGN.CENTER, bold=True)
card(s, Inches(1.6), Inches(1.0), Inches(0.03), Inches(1.4), BORDER, r=0)
txt(s, Inches(1.8), Inches(0.95), Inches(10.5), Inches(0.4), "课件智能生成", 22, WHITE, True)
txt(s, Inches(1.8), Inches(1.3), Inches(10.5), Inches(0.3), "精品课件自动化生成", 15, LIGHT)
bl = [("-","输入课程大纲，AI 自动生成课件",14,MUTED),("-","内置小龙虾三部曲教学框架",14,MUTED),("-","支持代码演示、公式渲染",14,MUTED),("-","2小时教案压缩至10分钟",14,LIGHT)]
bullets(s, Inches(1.8), Inches(1.62), Inches(10.5), Inches(0.9), bl)

# Scene 2
card(s, Inches(0.5), Inches(2.8), Inches(12.3), Inches(1.7), CARD, BORDER, 1)
txt(s, Inches(0.8), Inches(2.9), Inches(0.7), Inches(0.5), "02", 28, LIGHT, align=PP_ALIGN.CENTER, bold=True)
card(s, Inches(1.6), Inches(2.9), Inches(0.03), Inches(1.4), BORDER, r=0)
txt(s, Inches(1.8), Inches(2.85), Inches(10.5), Inches(0.4), "教学案例开发", 22, WHITE, True)
txt(s, Inches(1.8), Inches(3.2), Inches(10.5), Inches(0.3), "烟草数据挖掘 / 网络动画 / 微信小程序", 15, LIGHT)
bl = [("-","烟草数据挖掘：真实行业数据+完整流程",14,MUTED),("-","计算机网络16章教学动画",14,MUTED),("-","微信小程序：从0到1完整实战",14,MUTED),("-","基于阿里云智能体的教学实践",14,LIGHT)]
bullets(s, Inches(1.8), Inches(3.52), Inches(10.5), Inches(0.9), bl)

# Scene 3
card(s, Inches(0.5), Inches(4.7), Inches(12.3), Inches(1.7), CARD, BORDER, 1)
txt(s, Inches(0.8), Inches(4.8), Inches(0.7), Inches(0.5), "03", 28, LIGHT, align=PP_ALIGN.CENTER, bold=True)
card(s, Inches(1.6), Inches(4.8), Inches(0.03), Inches(1.4), BORDER, r=0)
txt(s, Inches(1.8), Inches(4.75), Inches(10.5), Inches(0.4), "论文协作写作", 22, WHITE, True)
txt(s, Inches(1.8), Inches(5.1), Inches(10.5), Inches(0.3), "从选题到 IEEE 成稿，全流程 AI 辅助", 15, LIGHT)
bl = [("-","选题：AI分析热点，推荐方向",14,MUTED),("-","写作：智能体辅助文献综述",14,MUTED),("-","格式：自动符合IEEE模板",14,MUTED),("-","审阅：多轮智能审校提升质量",14,LIGHT)]
bullets(s, Inches(1.8), Inches(5.42), Inches(10.5), Inches(0.9), bl)

# ============ SLIDE 4: SCENARIOS (DOWN) ============
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
card(s, Inches(0), Inches(0), Inches(13.333), Inches(0.05), LIGHT, r=0)
txt(s, Inches(1.2), Inches(0.35), Inches(11), Inches(0.4), "五大教学场景（下）", 16, LIGHT, True)

# Scene 4
card(s, Inches(0.5), Inches(0.9), Inches(12.3), Inches(1.5), CARD, BORDER, 1)
txt(s, Inches(0.8), Inches(1.0), Inches(0.7), Inches(0.5), "04", 28, LIGHT, align=PP_ALIGN.CENTER, bold=True)
card(s, Inches(1.6), Inches(1.0), Inches(0.03), Inches(1.2), BORDER, r=0)
txt(s, Inches(1.8), Inches(0.95), Inches(10.5), Inches(0.4), "教学视频制作", 22, WHITE, True)
txt(s, Inches(1.8), Inches(1.3), Inches(10.5), Inches(0.3), "PPT 自动转教学视频", 15, LIGHT)
bl = [("-","PPT课件一键转教学视频",14,MUTED),("-","智能语音合成自动生成配音",14,MUTED),("-","自动添加转场和字幕",14,LIGHT)]
bullets(s, Inches(1.8), Inches(1.6), Inches(10.5), Inches(0.75), bl)

# Scene 5
card(s, Inches(0.5), Inches(2.6), Inches(12.3), Inches(1.3), CARD, BORDER, 1)
txt(s, Inches(0.8), Inches(2.7), Inches(0.7), Inches(0.5), "05", 28, LIGHT, align=PP_ALIGN.CENTER, bold=True)
card(s, Inches(1.6), Inches(2.7), Inches(0.03), Inches(1.0), BORDER, r=0)
txt(s, Inches(1.8), Inches(2.65), Inches(10.5), Inches(0.4), "教案快速制作", 22, WHITE, True)
txt(s, Inches(1.8), Inches(3.0), Inches(10.5), Inches(0.3), "2小时教案 -> 10分钟完成", 15, LIGHT)
bl = [("-","输入目标，AI自动生成教案",14,MUTED),("-","自动匹配案例和练习题",14,LIGHT)]
bullets(s, Inches(1.8), Inches(3.25), Inches(10.5), Inches(0.55), bl)

# Metrics
card(s, Inches(0.5), Inches(4.2), Inches(12.3), Inches(2.5), RGBColor(0x0E,0x1E,0x3F), BLUE, 2)
txt(s, Inches(0.8), Inches(4.3), Inches(11.5), Inches(0.4), "核心数据", 20, LIGHT, True)
m = [("效率提升","20倍","教案从2h降至10min"),("案例覆盖","16章","计算机网络全套动画"),("实战项目","3+","数据挖掘/小程序/动画"),("视频转换","一键","PPT自动转视频")]
for i,(label,val,desc) in enumerate(m):
    col = i%2; row = i//2
    x = Inches(0.8)+col*Inches(5.9); y = Inches(4.8)+row*Inches(0.8)
    txt(s, x, y, Inches(5.5), Inches(0.25), label, 14, MUTED)
    txt(s, x, y+Inches(0.2), Inches(5.5), Inches(0.35), val, 22, WHITE, True)
    txt(s, x, y+Inches(0.5), Inches(5.5), Inches(0.22), desc, 13, MUTED)

# ============ SLIDE 5: AUTHOR ============
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
card(s, Inches(0), Inches(0), Inches(13.333), Inches(0.05), LIGHT, r=0)
txt(s, Inches(1.2), Inches(0.35), Inches(11), Inches(0.4), "主讲人", 16, LIGHT, True)

# Photo
photo = "/home/admin/.openclaw/workspace/poster_image1.png"
if os.path.exists(photo):
    try: s.shapes.add_picture(photo, Inches(0.8), Inches(1.0), Inches(2.6), Inches(3.3))
    except: pass
txt(s, Inches(0.8), Inches(4.4), Inches(2.6), Inches(0.4), "诸葛斌", 20, WHITE, align=PP_ALIGN.CENTER, bold=True)
txt(s, Inches(0.8), Inches(4.8), Inches(2.6), Inches(0.35), "教授", 14, MUTED, align=PP_ALIGN.CENTER)

txt(s, Inches(3.8), Inches(1.0), Inches(8.5), Inches(0.65), "诸葛斌 教授", 36, WHITE, True)
txt(s, Inches(3.8), Inches(1.7), Inches(8.5), Inches(0.35), "浙江工商大学 - 信息与电子工程学院", 16, BLUE)
txt(s, Inches(3.8), Inches(2.0), Inches(8.5), Inches(0.3), "萨塞克斯人工智能学院", 16, BLUE)
card(s, Inches(3.8), Inches(2.4), Inches(1.2), Inches(0.03), BLUE, r=0)

a = ["研究方向：互联网应用开发与AI教育","获浙江省技术发明一等奖","联合阿里钉钉撰写首本低代码教材","2025全国高校AI教育大会优秀案例","人工智能背景下智能体教学实践"]
for i,t in enumerate(a):
    y = Inches(2.6)+Inches(i*0.6)
    card(s, Inches(3.8), y, Inches(0.05), Inches(0.45), BLUE, r=0)
    txt(s, Inches(4.0), y, Inches(8), Inches(0.45), t, 16, MUTED)

# ============ SLIDE 6: BOOK ============
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
card(s, Inches(0), Inches(0), Inches(13.333), Inches(0.05), LIGHT, r=0)
txt(s, Inches(1.2), Inches(0.35), Inches(11), Inches(0.4), "所用图书", 16, LIGHT, True)
txt(s, Inches(1.2), Inches(0.85), Inches(11), Inches(0.6), "Manus 智能体全攻略", 32, WHITE, True)
txt(s, Inches(1.2), Inches(1.4), Inches(11), Inches(0.35), "清华大学出版社", 16, BLUE)

it = [("国内首本","智能体教学实战指南"),("覆盖全面","从基础到高级应用"),("实战导向","丰富案例和代码"),("高校适用","专为教师教学设计")]
for i,(l,d) in enumerate(it):
    x = Inches(0.7)+Inches(i*3.0)
    card(s, x, Inches(2.0), Inches(2.7), Inches(1.0), CARD, BORDER, 1)
    txt(s, x+Inches(0.2), Inches(2.1), Inches(2.3), Inches(0.35), l, 16, LIGHT, True)
    txt(s, x+Inches(0.2), Inches(2.45), Inches(2.3), Inches(0.35), d, 13, MUTED)

txt(s, Inches(1.2), Inches(3.3), Inches(11), Inches(0.4), "核心内容", 20, WHITE, True)
h = ["智能体基础概念与架构设计","Manus平台使用指南","教学场景应用案例","阿里云智能体教学实践","小龙虾AI协同工作流"]
bl = [("-",x,14,MUTED) for x in h]
bullets(s, Inches(1.2), Inches(3.7), Inches(11), Inches(1.1), bl)

# ============ SLIDE 7: BENEFITS ============
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
card(s, Inches(0), Inches(0), Inches(13.333), Inches(0.05), LIGHT, r=0)
txt(s, Inches(1.2), Inches(0.35), Inches(11), Inches(0.4), "入群福利", 16, LIGHT, True)
txt(s, Inches(1.2), Inches(0.85), Inches(11), Inches(0.5), "扫码加入读者服务群", 28, WHITE, True)
txt(s, Inches(1.2), Inches(1.3), Inches(11), Inches(0.3), "此群长期有效", 14, MUTED)

b = [("AI体验","小龙虾AI体验","群内部署OpenClaw，实时体验AI能力"),("资料包","教学资料包","课件/案例/动画/小程序全套资源"),("赠书","免费样书抽奖","直播间10本样书赠送"),("社区","教学交流社区","高校教师AI教学交流分享")]
for i,(tag,title,desc) in enumerate(b):
    col=i%2; row=i//2
    x=Inches(0.6)+col*Inches(6.1); y=Inches(1.7)+row*Inches(2.2)
    card(s, x, y, Inches(5.8), Inches(2.0), CARD, BORDER, 1)
    card(s, x+Inches(0.2), y+Inches(0.15), Inches(0.5), Inches(0.5), RGBColor(0x1E,0x3A,0x5F), BLUE, 1)
    txt(s, x+Inches(0.22), y+Inches(0.18), Inches(0.46), Inches(0.44), tag, 12, LIGHT, align=PP_ALIGN.CENTER, bold=True)
    txt(s, x+Inches(0.85), y+Inches(0.15), Inches(4.6), Inches(0.35), title, 18, WHITE, True)
    txt(s, x+Inches(0.85), y+Inches(0.55), Inches(4.6), Inches(1.3), desc, 13, MUTED, sp=1.25)

# ============ SLIDE 8: QR ============
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
card(s, Inches(0), Inches(0), Inches(13.333), Inches(0.05), LIGHT, r=0)
txt(s, Inches(1.2), Inches(0.35), Inches(11), Inches(0.4), "立即行动", 16, LIGHT, True)
txt(s, Inches(0.8), Inches(1.0), Inches(5), Inches(0.7), "6月27日 周五\n下午 15:00-16:00", 24, WHITE, True, sp=1.4)
txt(s, Inches(0.8), Inches(1.9), Inches(5), Inches(0.45), "主讲人：诸葛斌 教授", 20, BLUE, True)
txt(s, Inches(0.8), Inches(2.4), Inches(5), Inches(0.5), "基于 Manus 智能体全攻略\n清华大学出版社", 16, MUTED, sp=1.3)

qr = "/home/admin/.openclaw/workspace/lobster-network/qr_codes.png"
if os.path.exists(qr):
    try: s.shapes.add_picture(qr, Inches(6.2), Inches(0.6), Inches(6.2), Inches(4.2))
    except:
        card(s, Inches(6.2), Inches(0.6), Inches(6.2), Inches(4.2), CARD, BLUE, 2)
        txt(s, Inches(6.2), Inches(2.4), Inches(6.2), Inches(0.5), "扫码预约直播 + 进群", 20, WHITE, align=PP_ALIGN.CENTER, bold=True)
else:
    card(s, Inches(6.2), Inches(0.6), Inches(6.2), Inches(4.2), CARD, BLUE, 2)
    txt(s, Inches(6.2), Inches(2.4), Inches(6.2), Inches(0.5), "扫码预约直播 + 进群", 20, WHITE, align=PP_ALIGN.CENTER, bold=True)

card(s, Inches(0.7), Inches(5.2), Inches(11.9), Inches(1.2), RGBColor(0x0E,0x1E,0x3F), BLUE, 2)
txt(s, Inches(0.9), Inches(5.3), Inches(11.5), Inches(0.45), "立即扫码预约直播 + 加入读者群！", 20, LIGHT, align=PP_ALIGN.CENTER, bold=True)
txt(s, Inches(0.9), Inches(5.75), Inches(11.5), Inches(0.35), "群长期有效 - 教学资料包 - AI体验 - 样书抽奖", 14, MUTED, align=PP_ALIGN.CENTER)

# Save
out = "/home/admin/.openclaw/workspace/lobster-network/直播课程汇报_优化版.pptx"
prs.save(out)
print(f"Done: {out} ({len(prs.slides)} slides)")
