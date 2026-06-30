#!/usr/bin/env python3
"""
生成「智能体赋能高校教学新范式」直播课程汇报 PPTX - 布局优化版
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

DARK_BG     = RGBColor(0x0B, 0x11, 0x20)
BLUE_ACCENT = RGBColor(0x3B, 0x82, 0xF6)
LIGHT_BLUE  = RGBColor(0x60, 0xA5, 0xFA)
WHITE       = RGBColor(0xF0, 0xF4, 0xF8)
MUTED       = RGBColor(0x94, 0xA3, 0xB8)
CARD_BG     = RGBColor(0x1E, 0x29, 0x3B)
GREEN       = RGBColor(0x34, 0xD3, 0x99)

def add_bg(slide, color=DARK_BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_card(slide, left, top, width, height, fill_color=CARD_BG, border_color=None, border_width=None, radius_idx=0.1):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(border_width or 1)
    else:
        shape.line.fill.background()
    shape.adjustments[0] = radius_idx
    return shape

def add_text(slide, left, top, width, height, text, size=20, color=WHITE, bold=False, align=PP_ALIGN.LEFT, font_name="微软雅黑", line_spacing=1.3):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    p.line_spacing = line_spacing
    return tf

def add_bullet_list(slide, left, top, width, height, items):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (bullet, text, size, color) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet + " " + text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = "微软雅黑"
        p.space_after = Pt(8)
        p.line_spacing = 1.25
    return tf

# ═══ Slide 1: Cover ═══
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_card(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), LIGHT_BLUE, radius_idx=0)
add_text(slide, Inches(1.2), Inches(0.8), Inches(5), Inches(0.5), "直播课程 · 清华大学出版社", 18, LIGHT_BLUE, bold=True)
add_text(slide, Inches(0.8), Inches(2.0), Inches(11.7), Inches(1.8), "智能体赋能高校教学新范式", 50, WHITE, bold=True, align=PP_ALIGN.CENTER, line_spacing=1.2)
add_text(slide, Inches(1.2), Inches(3.6), Inches(11), Inches(0.7), "小龙虾 + Manus 一站式解决方案", 30, BLUE_ACCENT, bold=True, align=PP_ALIGN.CENTER)
add_card(slide, Inches(1.2), Inches(4.6), Inches(11), Inches(1.0), CARD_BG, RGBColor(0x33, 0x41, 0x55), 1)
add_text(slide, Inches(1.5), Inches(4.75), Inches(3.3), Inches(0.7), "诸葛斌 教授", 22, WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, Inches(5.0), Inches(4.75), Inches(3.3), Inches(0.7), "6月27日（周五）", 22, WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, Inches(8.5), Inches(4.75), Inches(3.3), Inches(0.7), "15:00 - 16:00", 22, WHITE, bold=True, align=PP_ALIGN.CENTER)
add_card(slide, Inches(3.5), Inches(6.0), Inches(6.3), Inches(0.6), RGBColor(0x1E, 0x3A, 0x5F), BLUE_ACCENT, 2)
add_text(slide, Inches(3.7), Inches(6.05), Inches(5.9), Inches(0.5), "基于 Manus 智能体全攻略", 18, WHITE, PP_ALIGN.CENTER)

# ═══ Slide 2: 痛点与解决方案 ═══
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_card(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), LIGHT_BLUE, radius_idx=0)
add_text(slide, Inches(1.2), Inches(0.4), Inches(11), Inches(0.5), "背景与痛点", 18, LIGHT_BLUE, bold=True)
add_text(slide, Inches(1.2), Inches(1.0), Inches(11), Inches(0.7), "高校教师面临的 AI 教学困境", 36, WHITE, bold=True)

pains = [
    ("课件制作耗时", "一份精品课件需要数小时甚至数天"),
    ("案例更新滞后", "教学案例跟不上技术发展速度"),
    ("论文写作困难", "从选题到成稿周期长、效率低"),
    ("视频制作复杂", "需要专业剪辑技能和软件"),
]
for i, (title, desc) in enumerate(pains):
    y = Inches(1.9) + Inches(i * 1.2)
    add_card(slide, Inches(0.6), y, Inches(5.6), Inches(1.05), CARD_BG, RGBColor(0xFF, 0x6B, 0x6B), 1)
    add_text(slide, Inches(0.9), y + Inches(0.12), Inches(5), Inches(0.4), title, 20, RGBColor(0xFF, 0x6B, 0x6B), bold=True)
    add_text(slide, Inches(0.9), y + Inches(0.52), Inches(5), Inches(0.4), desc, 16, MUTED)

add_text(slide, Inches(6.3), Inches(3.5), Inches(0.7), Inches(0.8), ">", 36, BLUE_ACCENT, PP_ALIGN.CENTER)

solutions = [
    ("课件生成", "2小时教案 -> 10分钟"),
    ("案例开发", "自动挖掘 + 可视化展示"),
    ("论文协作", "选题到成稿全流程"),
    ("视频制作", "PPT 自动转教学视频"),
]
for i, (title, desc) in enumerate(solutions):
    y = Inches(1.9) + Inches(i * 1.2)
    add_card(slide, Inches(7.0), y, Inches(5.6), Inches(1.05), RGBColor(0x0E, 0x3A, 0x2F), GREEN, 1)
    add_text(slide, Inches(7.3), y + Inches(0.12), Inches(5), Inches(0.4), title, 20, GREEN, bold=True)
    add_text(slide, Inches(7.3), y + Inches(0.52), Inches(5), Inches(0.4), desc, 16, MUTED)

# ═══ Slide 3: 五大场景（上） ═══
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_card(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), LIGHT_BLUE, radius_idx=0)
add_text(slide, Inches(1.2), Inches(0.4), Inches(11), Inches(0.5), "五大教学场景（上）", 18, LIGHT_BLUE, bold=True)

# Scene 1
add_card(slide, Inches(0.6), Inches(1.0), Inches(12.1), Inches(1.8), CARD_BG, RGBColor(0x33, 0x41, 0x55), 1)
add_text(slide, Inches(0.9), Inches(1.15), Inches(0.8), Inches(0.7), "S1", 40, WHITE, bold=True, align=PP_ALIGN.CENTER)
add_card(slide, Inches(1.8), Inches(1.15), Inches(0.04), Inches(1.5), RGBColor(0x33, 0x41, 0x55), radius_idx=0)
add_text(slide, Inches(2.0), Inches(1.1), Inches(10), Inches(0.45), "场景 01 - 课件智能生成", 24, WHITE, bold=True)
add_text(slide, Inches(2.0), Inches(1.5), Inches(10), Inches(0.35), "精品课件自动化生成", 17, LIGHT_BLUE)
bullet_items = [
    ("-", "输入课程大纲，AI 自动生成交互式课件", 15, MUTED),
    ("-", "内置小龙虾三部曲教学框架", 15, MUTED),
    ("-", "支持代码演示、公式渲染、动画交互", 15, MUTED),
    ("-", "2 小时教案制作压缩至 10 分钟", 15, LIGHT_BLUE),
]
add_bullet_list(slide, Inches(2.0), Inches(1.85), Inches(10), Inches(1.0), bullet_items)

# Scene 2
add_card(slide, Inches(0.6), Inches(3.0), Inches(12.1), Inches(1.8), CARD_BG, RGBColor(0x33, 0x41, 0x55), 1)
add_text(slide, Inches(0.9), Inches(3.15), Inches(0.8), Inches(0.7), "S2", 40, WHITE, bold=True, align=PP_ALIGN.CENTER)
add_card(slide, Inches(1.8), Inches(3.15), Inches(0.04), Inches(1.5), RGBColor(0x33, 0x41, 0x55), radius_idx=0)
add_text(slide, Inches(2.0), Inches(3.1), Inches(10), Inches(0.45), "场景 02 - 教学案例开发", 24, WHITE, bold=True)
add_text(slide, Inches(2.0), Inches(3.5), Inches(10), Inches(0.35), "烟草数据挖掘 / 网络课程动画 / 微信小程序", 17, LIGHT_BLUE)
bullet_items = [
    ("-", "烟草数据挖掘案例：真实行业数据 + 完整分析流程", 15, MUTED),
    ("-", "计算机网络 16 章教学动画：可视化复杂概念", 15, MUTED),
    ("-", "微信小程序开发案例：从 0 到 1 完整项目实战", 15, MUTED),
    ("-", "基于阿里云智能体的人工智能教学实践", 15, LIGHT_BLUE),
]
add_bullet_list(slide, Inches(2.0), Inches(3.85), Inches(10), Inches(1.0), bullet_items)

# Scene 3
add_card(slide, Inches(0.6), Inches(5.0), Inches(12.1), Inches(1.8), CARD_BG, RGBColor(0x33, 0x41, 0x55), 1)
add_text(slide, Inches(0.9), Inches(5.15), Inches(0.8), Inches(0.7), "S3", 40, WHITE, bold=True, align=PP_ALIGN.CENTER)
add_card(slide, Inches(1.8), Inches(5.15), Inches(0.04), Inches(1.5), RGBColor(0x33, 0x41, 0x55), radius_idx=0)
add_text(slide, Inches(2.0), Inches(5.1), Inches(10), Inches(0.45), "场景 03 - 论文协作写作", 24, WHITE, bold=True)
add_text(slide, Inches(2.0), Inches(5.5), Inches(10), Inches(0.35), "从选题到 IEEE 成稿，全流程 AI 辅助", 17, LIGHT_BLUE)
bullet_items = [
    ("-", "选题阶段：AI 分析研究热点，推荐创新方向", 15, MUTED),
    ("-", "写作阶段：智能体辅助文献综述、方法论述", 15, MUTED),
    ("-", "格式规范：自动符合 IEEE 模板要求", 15, MUTED),
    ("-", "协作审阅：多轮智能审校提升质量", 15, LIGHT_BLUE),
]
add_bullet_list(slide, Inches(2.0), Inches(5.85), Inches(10), Inches(1.0), bullet_items)

# ═══ Slide 4: 五大场景（下） ═══
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_card(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), LIGHT_BLUE, radius_idx=0)
add_text(slide, Inches(1.2), Inches(0.4), Inches(11), Inches(0.5), "五大教学场景（下）", 18, LIGHT_BLUE, bold=True)

# Scene 4
add_card(slide, Inches(0.6), Inches(1.0), Inches(12.1), Inches(1.6), CARD_BG, RGBColor(0x33, 0x41, 0x55), 1)
add_text(slide, Inches(0.9), Inches(1.12), Inches(0.8), Inches(0.7), "S4", 40, WHITE, bold=True, align=PP_ALIGN.CENTER)
add_card(slide, Inches(1.8), Inches(1.12), Inches(0.04), Inches(1.3), RGBColor(0x33, 0x41, 0x55), radius_idx=0)
add_text(slide, Inches(2.0), Inches(1.08), Inches(10), Inches(0.45), "场景 04 - 教学视频制作", 24, WHITE, bold=True)
add_text(slide, Inches(2.0), Inches(1.48), Inches(10), Inches(0.35), "PPT 自动转教学视频，无需专业剪辑", 17, LIGHT_BLUE)
bullet_items = [
    ("-", "PPT 课件一键转换为教学视频", 15, MUTED),
    ("-", "智能语音合成，自动生成讲解配音", 15, MUTED),
    ("-", "自动添加转场动画和字幕", 15, LIGHT_BLUE),
]
add_bullet_list(slide, Inches(2.0), Inches(1.82), Inches(10), Inches(0.7), bullet_items)

# Scene 5
add_card(slide, Inches(0.6), Inches(2.8), Inches(12.1), Inches(1.4), CARD_BG, RGBColor(0x33, 0x41, 0x55), 1)
add_text(slide, Inches(0.9), Inches(2.92), Inches(0.8), Inches(0.7), "S5", 40, WHITE, bold=True, align=PP_ALIGN.CENTER)
add_card(slide, Inches(1.8), Inches(2.92), Inches(0.04), Inches(1.1), RGBColor(0x33, 0x41, 0x55), radius_idx=0)
add_text(slide, Inches(2.0), Inches(2.88), Inches(10), Inches(0.45), "场景 05 - 教案快速制作", 24, WHITE, bold=True)
add_text(slide, Inches(2.0), Inches(3.28), Inches(10), Inches(0.35), "2 小时教案 -> 10 分钟完成", 17, LIGHT_BLUE)
bullet_items = [
    ("-", "输入教学目标和知识点，AI 自动生成教案", 15, MUTED),
    ("-", "自动匹配教学案例和练习题", 15, LIGHT_BLUE),
]
add_bullet_list(slide, Inches(2.0), Inches(3.58), Inches(10), Inches(0.6), bullet_items)

# Metrics
add_card(slide, Inches(0.6), Inches(4.5), Inches(12.1), Inches(2.3), RGBColor(0x0E, 0x1E, 0x3F), BLUE_ACCENT, 2)
add_text(slide, Inches(0.9), Inches(4.65), Inches(11.5), Inches(0.5), "核心数据", 22, LIGHT_BLUE, bold=True)
metrics = [
    ("效率提升", "20 倍", "教案从 2h 降至 10min"),
    ("案例覆盖", "16 章", "计算机网络全套动画"),
    ("实战项目", "3+ 个", "数据挖掘/小程序/动画"),
    ("视频转换", "一键", "PPT 自动转视频"),
]
for i, (label, value, desc) in enumerate(metrics):
    col = i % 2
    row = i // 2
    x = Inches(0.9) + col * Inches(5.9)
    y = Inches(5.15) + row * Inches(0.75)
    add_text(slide, x, y, Inches(5.5), Inches(0.3), label, 16, MUTED)
    add_text(slide, x, y + Inches(0.25), Inches(5.5), Inches(0.4), value, 24, WHITE, bold=True)
    add_text(slide, x, y + Inches(0.6), Inches(5.5), Inches(0.25), desc, 14, MUTED)

# ═══ Slide 5: 主讲人 ═══
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_card(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), LIGHT_BLUE, radius_idx=0)
add_text(slide, Inches(1.2), Inches(0.4), Inches(11), Inches(0.5), "主讲人", 18, LIGHT_BLUE, bold=True)

author_photo = "/home/admin/.openclaw/workspace/poster_image1.png"
if os.path.exists(author_photo):
    try:
        slide.shapes.add_picture(author_photo, Inches(1.0), Inches(1.2), Inches(2.8), Inches(3.5))
    except:
        pass
add_text(slide, Inches(1.0), Inches(4.8), Inches(2.8), Inches(0.5), "诸葛斌", 22, WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, Inches(1.0), Inches(5.3), Inches(2.8), Inches(0.4), "教授", 16, MUTED, PP_ALIGN.CENTER)

add_text(slide, Inches(4.2), Inches(1.2), Inches(8.5), Inches(0.7), "诸葛斌 教授", 40, WHITE, bold=True)
add_text(slide, Inches(4.2), Inches(1.95), Inches(8.5), Inches(0.4), "浙江工商大学 - 信息与电子工程学院", 18, BLUE_ACCENT)
add_text(slide, Inches(4.2), Inches(2.3), Inches(8.5), Inches(0.35), "萨塞克斯人工智能学院", 18, BLUE_ACCENT)
add_card(slide, Inches(4.2), Inches(2.75), Inches(1.5), Inches(0.04), BLUE_ACCENT, radius_idx=0)

achievements = [
    "研究方向：互联网应用开发与 AI 教育",
    "获浙江省技术发明一等奖",
    "联合阿里钉钉撰写国内首本低代码开发教材",
    "2025 全国高校 AI 教育大会优秀案例一等奖",
    "人工智能背景下基于阿里云的智能体教学实践",
]
for i, text in enumerate(achievements):
    y = Inches(3.0) + Inches(i * 0.65)
    add_card(slide, Inches(4.2), y, Inches(0.06), Inches(0.5), BLUE_ACCENT, radius_idx=0)
    add_text(slide, Inches(4.5), y, Inches(8), Inches(0.5), text, 18, MUTED)

# ═══ Slide 6: 图书介绍 ═══
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_card(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), LIGHT_BLUE, radius_idx=0)
add_text(slide, Inches(1.2), Inches(0.4), Inches(11), Inches(0.5), "所用图书", 18, LIGHT_BLUE, bold=True)
add_text(slide, Inches(1.2), Inches(1.0), Inches(11), Inches(0.7), "Manus 智能体全攻略", 36, WHITE, bold=True)
add_text(slide, Inches(1.2), Inches(1.65), Inches(11), Inches(0.4), "清华大学出版社", 18, BLUE_ACCENT)

info_items = [
    ("", "国内首本", "智能体教学实战指南"),
    ("", "覆盖全面", "从基础到高级应用"),
    ("", "实战导向", "丰富案例和代码"),
    ("", "高校适用", "专为教师教学设计"),
]
for i, (icon, label, desc) in enumerate(info_items):
    x = Inches(0.8) + Inches(i * 3.0)
    add_card(slide, x, Inches(2.3), Inches(2.7), Inches(1.1), CARD_BG, RGBColor(0x33, 0x41, 0x55), 1)
    add_text(slide, x + Inches(0.8), Inches(2.42), Inches(1.8), Inches(0.4), label, 18, LIGHT_BLUE, bold=True)
    add_text(slide, x + Inches(0.8), Inches(2.82), Inches(1.8), Inches(0.35), desc, 14, MUTED)

add_text(slide, Inches(1.2), Inches(3.7), Inches(11), Inches(0.5), "本书核心内容", 22, WHITE, bold=True)
highlights = [
    "智能体基础概念与架构设计",
    "Manus 平台使用指南与最佳实践",
    "教学场景应用案例",
    "阿里云智能体教学实践",
    "小龙虾 AI 助手协同工作流",
]
bullet_items = [("-", item, 16, MUTED) for item in highlights]
add_bullet_list(slide, Inches(1.2), Inches(4.2), Inches(11), Inches(1.2), bullet_items)

# ═══ Slide 7: 入群福利 ══
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_card(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), LIGHT_BLUE, radius_idx=0)
add_text(slide, Inches(1.2), Inches(0.4), Inches(11), Inches(0.5), "入群福利", 18, LIGHT_BLUE, bold=True)
add_text(slide, Inches(1.2), Inches(0.95), Inches(11), Inches(0.5), "扫码加入读者服务群", 32, WHITE, bold=True)
add_text(slide, Inches(1.2), Inches(1.4), Inches(11), Inches(0.35), "此群长期有效", 16, MUTED)

benefits = [
    ("AI", "小龙虾 AI 体验", "群内部署 OpenClaw 智能体，实时体验 AI 能力"),
    ("资料", "教学资料包", "课件 / 案例 / 动画 / 小程序全套资源"),
    ("赠书", "免费样书抽奖", "直播间专享 10 本样书赠送"),
    ("社区", "教学交流社区", "高校教师 AI 教学交流分享"),
]
for i, (tag, title, desc) in enumerate(benefits):
    col = i % 2
    row = i // 2
    x = Inches(0.7) + col * Inches(6.0)
    y = Inches(1.9) + row * Inches(2.2)
    add_card(slide, x, y, Inches(5.7), Inches(2.0), CARD_BG, RGBColor(0x33, 0x41, 0x55), 1)
    add_card(slide, x + Inches(0.2), y + Inches(0.15), Inches(0.5), Inches(0.5), RGBColor(0x1E, 0x3A, 0x5F), BLUE_ACCENT, 1)
    add_text(slide, x + Inches(0.25), y + Inches(0.2), Inches(0.4), Inches(0.4), tag, 14, LIGHT_BLUE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, x + Inches(0.9), y + Inches(0.15), Inches(4.5), Inches(0.4), title, 20, WHITE, bold=True)
    add_text(slide, x + Inches(0.9), y + Inches(0.65), Inches(4.5), Inches(1.2), desc, 14, MUTED, line_spacing=1.3)

# ═══ Slide 8: 预约 + 二维码 ═══
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_card(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), LIGHT_BLUE, radius_idx=0)
add_text(slide, Inches(1.2), Inches(0.4), Inches(11), Inches(0.5), "立即行动", 18, LIGHT_BLUE, bold=True)

add_text(slide, Inches(1.0), Inches(1.2), Inches(5), Inches(0.8), "6月27日（周五）\n下午 15:00 - 16:00", 26, WHITE, bold=True, line_spacing=1.5)
add_text(slide, Inches(1.0), Inches(2.2), Inches(5), Inches(0.5), "主讲人：诸葛斌 教授", 22, BLUE_ACCENT, bold=True)
add_text(slide, Inches(1.0), Inches(2.8), Inches(5), Inches(0.6), "基于 Manus 智能体全攻略\n清华大学出版社", 18, MUTED, line_spacing=1.4)

qr_path = "/home/admin/.openclaw/workspace/lobster-network/qr_codes.png"
if os.path.exists(qr_path):
    try:
        slide.shapes.add_picture(qr_path, Inches(6.5), Inches(0.8), Inches(5.8), Inches(4.0))
    except:
        add_card(slide, Inches(6.5), Inches(0.8), Inches(5.8), Inches(4.0), CARD_BG, BLUE_ACCENT, 2)
        add_text(slide, Inches(6.5), Inches(2.5), Inches(5.8), Inches(0.6), "扫码预约直播 + 进群", 22, WHITE, bold=True, align=PP_ALIGN.CENTER)
else:
    add_card(slide, Inches(6.5), Inches(0.8), Inches(5.8), Inches(4.0), CARD_BG, BLUE_ACCENT, 2)
    add_text(slide, Inches(6.5), Inches(2.5), Inches(5.8), Inches(0.6), "扫码预约直播 + 进群", 22, WHITE, bold=True, align=PP_ALIGN.CENTER)

add_card(slide, Inches(0.8), Inches(5.2), Inches(11.7), Inches(1.3), RGBColor(0x0E, 0x1E, 0x3F), BLUE_ACCENT, 2)
add_text(slide, Inches(1.0), Inches(5.35), Inches(11.3), Inches(0.5), "立即扫码预约直播 + 加入读者群！", 22, LIGHT_BLUE, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, Inches(1.0), Inches(5.85), Inches(11.3), Inches(0.4), "群长期有效 - 教学资料包 - AI 体验 - 样书抽奖", 16, MUTED, PP_ALIGN.CENTER)

# Save
output_path = "/home/admin/.openclaw/workspace/lobster-network/直播课程汇报_优化版.pptx"
prs.save(output_path)
print(f"PPTX saved: {output_path}")
print(f"Slides: {len(prs.slides)}")
