#!/usr/bin/env python3
"""
生成「智能体赋能高校教学新范式」直播课程汇报 PPTX — 完整版
包含：作者照片、二维码、详细案例内容
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── 配色方案 ──
DARK_BG     = RGBColor(0x0B, 0x11, 0x20)
BLUE_ACCENT = RGBColor(0x3B, 0x82, 0xF6)
LIGHT_BLUE  = RGBColor(0x60, 0xA5, 0xFA)
WHITE       = RGBColor(0xF0, 0xF4, 0xF8)
MUTED       = RGBColor(0x94, 0xA3, 0xB8)
CARD_BG     = RGBColor(0x1E, 0x29, 0x3B)
PURPLE      = RGBColor(0x7C, 0x3A, 0xED)

def add_bg(slide, color=DARK_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_card(slide, left, top, width, height, fill_color=CARD_BG, border_color=None, border_width=None, radius_idx=0.1):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(border_width or 1)
    else:
        shape.line.fill.background()
    shape.adjustments[0] = radius_idx
    return shape

def set_text(tf, text, size=24, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="微软雅黑"):
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment

def add_icon_label(slide, left, top, icon_text, label, label_color=MUTED, icon_size=32):
    txBox = slide.shapes.add_textbox(left, top, Inches(0.6), Inches(0.6))
    set_text(txBox.text_frame, icon_text, icon_size, WHITE)
    txBox = slide.shapes.add_textbox(left + Inches(0.8), top + Inches(0.1), Inches(4), Inches(0.5))
    set_text(txBox.text_frame, label, 18, label_color)

# ═══════════════════════════════════════════
# Slide 1: Cover
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

# Top accent gradient bar
add_card(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), LIGHT_BLUE, radius_idx=0)

# "直播课程" tag
txBox = slide.shapes.add_textbox(Inches(1.5), Inches(1.0), Inches(3), Inches(0.5))
set_text(txBox.text_frame, "📺 直播课程 · 清华大学出版社", 18, LIGHT_BLUE, bold=True)

# Main title
txBox = slide.shapes.add_textbox(Inches(1.2), Inches(1.8), Inches(11), Inches(1.6))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "智能体赋能高校教学新范式"
p.font.size = Pt(52)
p.font.color.rgb = WHITE
p.font.bold = True
p.font.name = "微软雅黑"
p.alignment = PP_ALIGN.LEFT

# Subtitle
txBox = slide.shapes.add_textbox(Inches(1.2), Inches(3.4), Inches(11), Inches(0.8))
set_text(txBox.text_frame, "小龙虾 + Manus 一站式解决方案", 32, BLUE_ACCENT, bold=True)

# Info bar
add_card(slide, Inches(1.2), Inches(4.6), Inches(11), Inches(1.2), CARD_BG, RGBColor(0x33, 0x41, 0x55), 1)
txBox = slide.shapes.add_textbox(Inches(1.5), Inches(4.8), Inches(3.5), Inches(0.8))
set_text(txBox.text_frame, " 诸葛斌 教授", 24, WHITE, bold=True)
txBox = slide.shapes.add_textbox(Inches(5.2), Inches(4.8), Inches(3.5), Inches(0.8))
set_text(txBox.text_frame, "📅 6月27日（周五）", 24, WHITE, bold=True)
txBox = slide.shapes.add_textbox(Inches(8.9), Inches(4.8), Inches(3), Inches(0.8))
set_text(txBox.text_frame, "⏰ 15:00 - 16:00", 24, WHITE, bold=True)

# Book badge
add_card(slide, Inches(1.2), Inches(6.1), Inches(5), Inches(0.7), RGBColor(0x1E, 0x3A, 0x5F), BLUE_ACCENT, 2)
txBox = slide.shapes.add_textbox(Inches(1.4), Inches(6.15), Inches(4.6), Inches(0.6))
set_text(txBox.text_frame, "📘 基于《Manus智能体全攻略》（清华大学出版社）", 20, WHITE)

# ═══════════════════════════════════════════
# Slide 2: 痛点与解决方案
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_card(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), LIGHT_BLUE, radius_idx=0)

txBox = slide.shapes.add_textbox(Inches(1.2), Inches(0.4), Inches(11), Inches(0.5))
set_text(txBox.text_frame, "背 景 与 痛 点", 18, LIGHT_BLUE, bold=True)

txBox = slide.shapes.add_textbox(Inches(1.2), Inches(1.0), Inches(11), Inches(0.9))
set_text(txBox.text_frame, "高校教师面临的 AI 教学困境", 40, WHITE, bold=True)

# Pain points (left)
pains = [
    ("😫", "课件制作耗时", "一份精品课件需要数小时甚至数天"),
    ("😰", "案例更新滞后", "教学案例跟不上技术发展速度"),
    ("😩", "论文写作困难", "从选题到成稿周期长、效率低"),
    ("", "视频制作复杂", "需要专业剪辑技能和软件"),
]
for i, (icon, title, desc) in enumerate(pains):
    y = Inches(2.1) + Inches(i * 0.85)
    add_card(slide, Inches(0.8), y, Inches(5.6), Inches(0.7), CARD_BG, RGBColor(0xFF, 0x6B, 0x6B), 1)
    txBox = slide.shapes.add_textbox(Inches(1.0), y + Inches(0.08), Inches(0.5), Inches(0.5))
    set_text(txBox.text_frame, icon, 28, WHITE)
    txBox = slide.shapes.add_textbox(Inches(1.5), y + Inches(0.05), Inches(4.7), Inches(0.35))
    set_text(txBox.text_frame, title, 22, RGBColor(0xFF, 0x6B, 0x6B), bold=True)
    txBox = slide.shapes.add_textbox(Inches(1.5), y + Inches(0.4), Inches(4.7), Inches(0.3))
    set_text(txBox.text_frame, desc, 17, MUTED)

# Arrow
txBox = slide.shapes.add_textbox(Inches(6.4), Inches(3.5), Inches(0.5), Inches(1))
set_text(txBox.text_frame, "⬇️", 36, BLUE_ACCENT, alignment=PP_ALIGN.CENTER)

# Solution (right)
txBox = slide.shapes.add_textbox(Inches(6.8), Inches(2.1), Inches(6), Inches(0.7))
set_text(txBox.text_frame, " 小龙虾 + Manus 解决方案", 32, LIGHT_BLUE, bold=True)

solutions = [
    ("✅", "课件生成", "2小时教案 → 10分钟"),
    ("✅", "案例开发", "自动挖掘 + 可视化展示"),
    ("✅", "论文协作", "选题→成稿全流程"),
    ("✅", "视频制作", "PPT 自动转教学视频"),
]
for i, (icon, title, desc) in enumerate(solutions):
    y = Inches(3.1) + Inches(i * 0.85)
    add_card(slide, Inches(6.8), y, Inches(5.6), Inches(0.7), RGBColor(0x0E, 0x3A, 0x2F), RGBColor(0x34, 0xD3, 0x99), 1)
    txBox = slide.shapes.add_textbox(Inches(7.0), y + Inches(0.08), Inches(0.4), Inches(0.5))
    set_text(txBox.text_frame, icon, 28, WHITE)
    txBox = slide.shapes.add_textbox(Inches(7.5), y + Inches(0.05), Inches(4.7), Inches(0.35))
    set_text(txBox.text_frame, title, 22, RGBColor(0x34, 0xD3, 0x99), bold=True)
    txBox = slide.shapes.add_textbox(Inches(7.5), y + Inches(0.4), Inches(4.7), Inches(0.3))
    set_text(txBox.text_frame, desc, 17, MUTED)

# ═══════════════════════════════════════════
# Slide 3: 五大场景详述（上）
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_card(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), LIGHT_BLUE, radius_idx=0)

txBox = slide.shapes.add_textbox(Inches(1.2), Inches(0.4), Inches(11), Inches(0.5))
set_text(txBox.text_frame, "五 大 教 学 场 景（上）", 18, LIGHT_BLUE, bold=True)

# Scene 1: 课件智能生成
add_card(slide, Inches(0.8), Inches(1.1), Inches(11.5), Inches(1.7), CARD_BG, RGBColor(0x33, 0x41, 0x55), 1)

add_icon_label(slide, Inches(1.0), Inches(1.2), "📚", "场景 01", BLUE_ACCENT)
txBox = slide.shapes.add_textbox(Inches(2.0), Inches(1.15), Inches(5), Inches(0.5))
set_text(txBox.text_frame, "课件智能生成", 28, WHITE, bold=True)
txBox = slide.shapes.add_textbox(Inches(2.0), Inches(1.55), Inches(10), Inches(0.4))
set_text(txBox.text_frame, '"小龙虾三部曲"精品课件自动化生成', 18, LIGHT_BLUE)

txBox = slide.shapes.add_textbox(Inches(1.0), Inches(1.85), Inches(11), Inches(0.9))
details = [
    "• 输入课程大纲，AI 自动生成交互式课件",
    "• 内置小龙虾三部曲教学框架（导入→探究→总结）",
    "• 支持代码演示、公式渲染、动画交互",
    "• 传统 2 小时教案制作 → 压缩至 10 分钟"
]
for i, line in enumerate(details):
    p = txBox.text_frame.add_paragraph() if i > 0 else txBox.text_frame.paragraphs[0]
    p.text = line
    p.font.size = Pt(16)
    p.font.color.rgb = MUTED
    p.font.name = "微软雅黑"

# Scene 2: 教学案例开发
add_card(slide, Inches(0.8), Inches(3.0), Inches(11.5), Inches(2.0), CARD_BG, RGBColor(0x33, 0x41, 0x55), 1)

add_icon_label(slide, Inches(1.0), Inches(3.1), "🔬", "场景 02", BLUE_ACCENT)
txBox = slide.shapes.add_textbox(Inches(2.0), Inches(3.05), Inches(5), Inches(0.5))
set_text(txBox.text_frame, "教学案例开发", 28, WHITE, bold=True)
txBox = slide.shapes.add_textbox(Inches(2.0), Inches(3.45), Inches(10), Inches(0.4))
set_text(txBox.text_frame, "烟草数据挖掘 · 网络课程动画 · 微信小程序实战", 18, LIGHT_BLUE)

txBox = slide.shapes.add_textbox(Inches(1.0), Inches(3.8), Inches(11), Inches(1.1))
details = [
    "• 烟草数据挖掘案例：真实行业数据 + 完整分析流程",
    "• 计算机网络 16 章教学动画：可视化复杂概念（TCP/IP、路由协议等）",
    "• 微信小程序开发案例：从 0 到 1 完整项目实战",
    "• 基于阿里云智能体的人工智能加教学实践"
]
for i, line in enumerate(details):
    p = txBox.text_frame.add_paragraph() if i > 0 else txBox.text_frame.paragraphs[0]
    p.text = line
    p.font.size = Pt(16)
    p.font.color.rgb = MUTED
    p.font.name = "微软雅黑"

# Scene 3: 论文协作写作
add_card(slide, Inches(0.8), Inches(5.2), Inches(11.5), Inches(1.8), CARD_BG, RGBColor(0x33, 0x41, 0x55), 1)

add_icon_label(slide, Inches(1.0), Inches(5.3), "📝", "场景 03", BLUE_ACCENT)
txBox = slide.shapes.add_textbox(Inches(2.0), Inches(5.25), Inches(5), Inches(0.5))
set_text(txBox.text_frame, "论文协作写作", 28, WHITE, bold=True)
txBox = slide.shapes.add_textbox(Inches(2.0), Inches(5.65), Inches(10), Inches(0.4))
set_text(txBox.text_frame, '从选题到 IEEE 成稿，全流程 AI 辅助', 18, LIGHT_BLUE)

txBox = slide.shapes.add_textbox(Inches(1.0), Inches(5.95), Inches(11), Inches(0.9))
details = [
    "• 选题阶段：AI 分析研究热点，推荐创新方向",
    "• 写作阶段：智能体辅助文献综述、方法论述、实验设计",
    "• 格式规范：自动符合 IEEE 模板要求，减少格式返工",
    "• 协作审阅：多轮智能审校，提升论文质量"
]
for i, line in enumerate(details):
    p = txBox.text_frame.add_paragraph() if i > 0 else txBox.text_frame.paragraphs[0]
    p.text = line
    p.font.size = Pt(16)
    p.font.color.rgb = MUTED
    p.font.name = "微软雅黑"

# ═══════════════════════════════════════════
# Slide 4: 五大场景详述（下）
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_card(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), LIGHT_BLUE, radius_idx=0)

txBox = slide.shapes.add_textbox(Inches(1.2), Inches(0.4), Inches(11), Inches(0.5))
set_text(txBox.text_frame, "五 大 教 学 场 景（下）", 18, LIGHT_BLUE, bold=True)

# Scene 4: 教学视频制作
add_card(slide, Inches(0.8), Inches(1.1), Inches(11.5), Inches(1.7), CARD_BG, RGBColor(0x33, 0x41, 0x55), 1)

add_icon_label(slide, Inches(1.0), Inches(1.2), "🎬", "场景 04", BLUE_ACCENT)
txBox = slide.shapes.add_textbox(Inches(2.0), Inches(1.15), Inches(5), Inches(0.5))
set_text(txBox.text_frame, "教学视频制作", 28, WHITE, bold=True)
txBox = slide.shapes.add_textbox(Inches(2.0), Inches(1.55), Inches(10), Inches(0.4))
set_text(txBox.text_frame, "PPT 自动转教学视频，无需专业剪辑", 18, LIGHT_BLUE)

txBox = slide.shapes.add_textbox(Inches(1.0), Inches(1.85), Inches(11), Inches(0.9))
details = [
    "• PPT 课件一键转换为教学视频",
    "• 智能语音合成，自动生成讲解配音",
    "• 自动添加转场动画和字幕",
    "• 支持多平台格式导出（MP4、WebM 等）"
]
for i, line in enumerate(details):
    p = txBox.text_frame.add_paragraph() if i > 0 else txBox.text_frame.paragraphs[0]
    p.text = line
    p.font.size = Pt(16)
    p.font.color.rgb = MUTED
    p.font.name = "微软雅黑"

# Scene 5: 教案快速制作
add_card(slide, Inches(0.8), Inches(3.0), Inches(11.5), Inches(1.7), CARD_BG, RGBColor(0x33, 0x41, 0x55), 1)

add_icon_label(slide, Inches(1.0), Inches(3.1), "⚡", "场景 05", BLUE_ACCENT)
txBox = slide.shapes.add_textbox(Inches(2.0), Inches(3.05), Inches(5), Inches(0.5))
set_text(txBox.text_frame, "教案快速制作", 28, WHITE, bold=True)
txBox = slide.shapes.add_textbox(Inches(2.0), Inches(3.45), Inches(10), Inches(0.4))
set_text(txBox.text_frame, "2 小时教案制作 → 10 分钟完成", 18, LIGHT_BLUE)

txBox = slide.shapes.add_textbox(Inches(1.0), Inches(3.8), Inches(11), Inches(0.8))
details = [
    "• 输入教学目标和知识点，AI 自动生成完整教案",
    "• 自动匹配教学案例和练习题",
    "• 支持多版本教案生成（基础版/进阶版）"
]
for i, line in enumerate(details):
    p = txBox.text_frame.add_paragraph() if i > 0 else txBox.text_frame.paragraphs[0]
    p.text = line
    p.font.size = Pt(16)
    p.font.color.rgb = MUTED
    p.font.name = "微软雅黑"

# Key metrics box
add_card(slide, Inches(0.8), Inches(5.0), Inches(11.5), Inches(1.8), RGBColor(0x0E, 0x1E, 0x3F), BLUE_ACCENT, 2)
txBox = slide.shapes.add_textbox(Inches(1.0), Inches(5.15), Inches(11), Inches(0.5))
set_text(txBox.text_frame, " 核心数据", 22, LIGHT_BLUE, bold=True)

metrics = [
    ("⏱️ 效率提升", "20 倍", "教案制作时间从 2 小时降至 10 分钟"),
    ("📚 案例数量", "16 章", "计算机网络全套教学动画覆盖"),
    ("🎯 实战项目", "3+", "烟草数据挖掘、微信小程序、网络动画"),
]
for i, (label, value, desc) in enumerate(metrics):
    x = Inches(1.2) + Inches(i * 3.9)
    txBox = slide.shapes.add_textbox(x, Inches(5.7), Inches(3.5), Inches(0.4))
    set_text(txBox.text_frame, label, 16, MUTED)
    txBox = slide.shapes.add_textbox(x, Inches(6.1), Inches(3.5), Inches(0.5))
    set_text(txBox.text_frame, value, 28, WHITE, bold=True)
    txBox = slide.shapes.add_textbox(x, Inches(6.6), Inches(3.5), Inches(0.3))
    set_text(txBox.text_frame, desc, 14, MUTED)

# ═══════════════════════════════════════════
# Slide 5: 主讲人介绍
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_card(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), LIGHT_BLUE, radius_idx=0)

txBox = slide.shapes.add_textbox(Inches(1.2), Inches(0.4), Inches(11), Inches(0.5))
set_text(txBox.text_frame, "主 讲 人", 18, LIGHT_BLUE, bold=True)

# Author photo placeholder (left)
add_card(slide, Inches(1.0), Inches(1.2), Inches(3.0), Inches(4.0), RGBColor(0x1E, 0x3A, 0x5F), BLUE_ACCENT, 2)

# Try to add author photo
author_photo_path = "/home/admin/.openclaw/workspace/poster_image1.png"
if os.path.exists(author_photo_path):
    try:
        slide.shapes.add_picture(author_photo_path, Inches(1.05), Inches(1.25), Inches(2.9), Inches(3.9))
    except Exception as e:
        print(f"Photo add failed: {e}")

txBox = slide.shapes.add_textbox(Inches(1.0), Inches(5.3), Inches(3.0), Inches(0.4))
set_text(txBox.text_frame, "诸葛斌", 24, WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# Author info (right)
txBox = slide.shapes.add_textbox(Inches(4.5), Inches(1.2), Inches(8), Inches(0.8))
set_text(txBox.text_frame, "诸葛斌 教授", 44, WHITE, bold=True)

txBox = slide.shapes.add_textbox(Inches(4.5), Inches(2.0), Inches(8), Inches(0.5))
set_text(txBox.text_frame, "浙江工商大学 · 信息与电子工程学院 / 萨塞克斯人工智能学院", 20, BLUE_ACCENT)

# Divider
add_card(slide, Inches(4.5), Inches(2.6), Inches(2), Inches(0.04), BLUE_ACCENT, radius_idx=0)

# Achievements
achievements = [
    ("🏆", "浙江省技术发明一等奖"),
    ("", "联合阿里钉钉撰写国内首本低代码开发教材"),
    ("🎖️", "2025 全国高校人工智能教育大会优秀案例一等奖"),
    ("", '"人工智能+"背景下基于阿里云的智能体教学实践'),
    ("", "研究方向：互联网应用开发与 AI 教育"),
]
for i, (icon, text) in enumerate(achievements):
    y = Inches(2.9) + Inches(i * 0.55)
    txBox = slide.shapes.add_textbox(Inches(4.5), y, Inches(0.4), Inches(0.5))
    set_text(txBox.text_frame, icon, 22, WHITE)
    txBox = slide.shapes.add_textbox(Inches(5.0), y, Inches(7), Inches(0.5))
    set_text(txBox.text_frame, text, 20, MUTED)

# ══════════════════════════════════════════
# Slide 6: 图书介绍
# ══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_card(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), LIGHT_BLUE, radius_idx=0)

txBox = slide.shapes.add_textbox(Inches(1.2), Inches(0.4), Inches(11), Inches(0.5))
set_text(txBox.text_frame, "所 用 图 书", 18, LIGHT_BLUE, bold=True)

txBox = slide.shapes.add_textbox(Inches(1.2), Inches(1.0), Inches(11), Inches(0.8))
set_text(txBox.text_frame, "《Manus智能体全攻略》", 40, WHITE, bold=True)
txBox = slide.shapes.add_textbox(Inches(1.2), Inches(1.7), Inches(11), Inches(0.4))
set_text(txBox.text_frame, "清华大学出版社出版", 22, BLUE_ACCENT)

# Book info cards
info_cards = [
    ("📘", "国内首本", "智能体教学实战指南"),
    ("🎯", "覆盖全面", "从基础概念到高级应用"),
    ("", "实战导向", "丰富的教学案例和代码"),
    ("", "高校适用", "专为高校教师教学设计"),
]
for i, (icon, label, desc) in enumerate(info_cards):
    x = Inches(0.8) + Inches(i * 3.0)
    add_card(slide, x, Inches(2.4), Inches(2.8), Inches(1.2), CARD_BG, RGBColor(0x33, 0x41, 0x55), 1)
    txBox = slide.shapes.add_textbox(x + Inches(0.2), Inches(2.5), Inches(0.5), Inches(0.5))
    set_text(txBox.text_frame, icon, 32, WHITE)
    txBox = slide.shapes.add_textbox(x + Inches(0.8), Inches(2.5), Inches(2), Inches(0.4))
    set_text(txBox.text_frame, label, 20, LIGHT_BLUE, bold=True)
    txBox = slide.shapes.add_textbox(x + Inches(0.8), Inches(2.95), Inches(2), Inches(0.4))
    set_text(txBox.text_frame, desc, 16, MUTED)

# Book highlights
txBox = slide.shapes.add_textbox(Inches(1.2), Inches(3.9), Inches(11), Inches(0.6))
set_text(txBox.text_frame, "📋 本书核心内容", 24, WHITE, bold=True)

highlights = [
    "• 智能体基础概念与架构设计",
    "• Manus 平台使用指南与最佳实践",
    "• 教学场景应用案例（课件、案例、论文、视频）",
    "• 阿里云智能体教学实践与人才培养创新",
    "• 小龙虾 AI 助手与 Manus 的协同工作流"
]
txBox = slide.shapes.add_textbox(Inches(1.2), Inches(4.4), Inches(11), Inches(1.2))
for i, line in enumerate(highlights):
    p = txBox.text_frame.add_paragraph() if i > 0 else txBox.text_frame.paragraphs[0]
    p.text = line
    p.font.size = Pt(18)
    p.font.color.rgb = MUTED
    p.font.name = "微软雅黑"

# ═══════════════════════════════════════════
# Slide 7: 入群福利
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_card(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), LIGHT_BLUE, radius_idx=0)

txBox = slide.shapes.add_textbox(Inches(1.2), Inches(0.4), Inches(11), Inches(0.5))
set_text(txBox.text_frame, "入 群 福 利", 18, LIGHT_BLUE, bold=True)

txBox = slide.shapes.add_textbox(Inches(1.2), Inches(1.0), Inches(11), Inches(0.6))
set_text(txBox.text_frame, "扫码加入读者服务群", 36, WHITE, bold=True)
txBox = slide.shapes.add_textbox(Inches(1.2), Inches(1.55), Inches(11), Inches(0.4))
set_text(txBox.text_frame, "此群长期有效，群成员可获得以下专属福利", 18, MUTED)

# 4 benefit cards
benefits = [
    ("", "小龙虾 AI 体验", "群内部署 OpenClaw 智能体，实时体验课件生成、教案制作等 AI 能力，让 AI 助手帮你备课"),
    ("📦", "教学资料包", "小龙虾三部曲课件、烟草数据挖掘案例、16 章计算机网络教学动画、微信小程序案例等全套 Manus 实战案例"),
    ("🎁", "免费样书抽奖", "直播间专享 10 本《Manus智能体全攻略》免费样书赠送（名额有限，抽奖获得）"),
    ("💬", "教学交流社区", '智能体系列课程教学交流社区——高校教师 AI 教学实践交流、问题解答、经验分享'),
]
for i, (icon, title, desc) in enumerate(benefits):
    col = i % 2
    row = i // 2
    x = Inches(0.8) + col * Inches(6.0)
    y = Inches(2.1) + row * Inches(2.1)
    add_card(slide, x, y, Inches(5.7), Inches(1.9), CARD_BG, RGBColor(0x33, 0x41, 0x55), 1)
    txBox = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.15), Inches(0.6), Inches(0.6))
    set_text(txBox.text_frame, icon, 36, WHITE)
    txBox = slide.shapes.add_textbox(x + Inches(0.9), y + Inches(0.15), Inches(4.5), Inches(0.4))
    set_text(txBox.text_frame, title, 22, WHITE, bold=True)
    txBox = slide.shapes.add_textbox(x + Inches(0.9), y + Inches(0.65), Inches(4.5), Inches(1.1))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = desc
    p.font.size = Pt(15)
    p.font.color.rgb = MUTED
    p.font.name = "微软雅黑"

# ═══════════════════════════════════════════
# Slide 8: 预约直播 + 二维码
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_card(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), LIGHT_BLUE, radius_idx=0)

txBox = slide.shapes.add_textbox(Inches(1.2), Inches(0.4), Inches(11), Inches(0.5))
set_text(txBox.text_frame, "立 即 行 动", 18, LIGHT_BLUE, bold=True)

# Left: Event info
txBox = slide.shapes.add_textbox(Inches(1.0), Inches(1.2), Inches(5.5), Inches(1.0))
set_text(txBox.text_frame, "📅 6月27日（周五）\n⏰ 下午 15:00 - 16:00", 28, WHITE, bold=True)

txBox = slide.shapes.add_textbox(Inches(1.0), Inches(2.4), Inches(5.5), Inches(0.7))
set_text(txBox.text_frame, " 主讲人：诸葛斌 教授", 24, BLUE_ACCENT, bold=True)

txBox = slide.shapes.add_textbox(Inches(1.0), Inches(3.1), Inches(5.5), Inches(0.7))
set_text(txBox.text_frame, "📘 基于《Manus智能体全攻略》", 24, BLUE_ACCENT, bold=True)

txBox = slide.shapes.add_textbox(Inches(1.0), Inches(3.8), Inches(5.5), Inches(0.5))
set_text(txBox.text_frame, "清华大学出版社", 20, MUTED)

# Right: QR codes
try:
    qr_path = "/home/admin/.openclaw/workspace/lobster-network/qr_codes.png"
    if os.path.exists(qr_path):
        slide.shapes.add_picture(qr_path, Inches(6.8), Inches(1.0), Inches(5.5), Inches(3.5))
except Exception as e:
    print(f"QR code add failed: {e}")
    # Fallback: add text placeholders
    add_card(slide, Inches(6.8), Inches(1.0), Inches(5.5), Inches(3.5), CARD_BG, BLUE_ACCENT, 2)
    txBox = slide.shapes.add_textbox(Inches(6.8), Inches(2.0), Inches(5.5), Inches(0.8))
    set_text(txBox.text_frame, "📱 扫码预约直播\n📱 扫码进群享福利", 24, WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# Bottom CTA bar
add_card(slide, Inches(1.0), Inches(5.2), Inches(11.3), Inches(1.3), RGBColor(0x0E, 0x1E, 0x3F), BLUE_ACCENT, 2)
txBox = slide.shapes.add_textbox(Inches(1.2), Inches(5.35), Inches(11), Inches(0.5))
set_text(txBox.text_frame, "🔥 立即扫码预约直播 + 加入读者群，享 4 大专属福利！", 24, LIGHT_BLUE, bold=True, alignment=PP_ALIGN.CENTER)
txBox = slide.shapes.add_textbox(Inches(1.2), Inches(5.9), Inches(11), Inches(0.4))
set_text(txBox.text_frame, "群长期有效 · 教学资料包 · AI 体验 · 样书抽奖", 18, MUTED, alignment=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════
# Save
# ═══════════════════════════════════════════
output_path = "/home/admin/.openclaw/workspace/lobster-network/直播课程汇报_完整版.pptx"
prs.save(output_path)
print(f"✅ PPTX saved: {output_path}")
print(f"   Slides: {len(prs.slides)}")
