#!/usr/bin/env python3
"""生成南湖区中小学幼儿园后勤管理干部培训PPT"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ==================== 配色方案 ====================
# 南湖红 + 教育蓝 主题
DARK_RED = RGBColor(0x8B, 0x00, 0x00)       # 深红（南湖红）
RED = RGBColor(0xC4, 0x1A, 0x1A)             # 红
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
MED_GRAY = RGBColor(0x66, 0x66, 0x66)
BLUE = RGBColor(0x1E, 0x88, 0xE5)            # 教育蓝
DARK_BLUE = RGBColor(0x0D, 0x47, 0xA1)
GOLD = RGBColor(0xFF, 0xA0, 0x00)
LIGHT_BLUE = RGBColor(0xE3, 0xF2, 0xFD)
ACCENT_GREEN = RGBColor(0x2E, 0x7D, 0x32)
ACCENT_ORANGE = RGBColor(0xE6, 0x51, 0x00)

# ==================== 辅助函数 ====================

def add_bg_rect(slide, x, y, w, h, color, alpha=None):
    """添加背景矩形"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if alpha is not None:
        from lxml import etree
        ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
        srgb = shape.fill._fill.find(f'{{{ns}}}solidFill/{{{ns}}}srgbClr')
        if srgb is not None:
            alpha_elem = etree.SubElement(srgb, f'{{{ns}}}alpha')
            alpha_elem.set('val', str(int(alpha * 1000)))
    return shape

def add_text_box(slide, x, y, w, h, text, font_size=18, color=DARK_GRAY,
                 bold=False, alignment=PP_ALIGN.LEFT, font_name='微软雅黑',
                 line_spacing=1.3):
    """添加文本框"""
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(font_size * 0.3)
    if line_spacing != 1.0:
        p.line_spacing = Pt(font_size * line_spacing)
    return txBox

def add_bullet_list(slide, x, y, w, h, items, font_size=16, color=DARK_GRAY,
                    bullet_color=BLUE, font_name='微软雅黑', line_spacing=1.6,
                    bold_items=None):
    """添加项目符号列表"""
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = Pt(4)
        p.level = 0
        p.line_spacing = Pt(font_size * line_spacing)
        if bold_items and i in bold_items:
            p.font.bold = True
        # 项目符号
        from pptx.oxml.ns import qn
        pPr = p._p.get_or_add_pPr()
        buChar = pPr.makeelement(qn('a:buChar'), {'char': '●'})
        # 先删除已有的
        for child in list(pPr):
            if child.tag.endswith('buChar') or child.tag.endswith('buNone'):
                pPr.remove(child)
        pPr.append(buChar)
    return txBox

def add_numbered_list(slide, x, y, w, h, items, font_size=16, color=DARK_GRAY,
                      font_name='微软雅黑', line_spacing=1.6):
    """添加编号列表"""
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"{i+1}. {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = Pt(6)
        p.line_spacing = Pt(font_size * line_spacing)
    return txBox

def add_card(slide, x, y, w, h, title, content_lines, title_color=WHITE,
             bg_color=BLUE, content_color=WHITE, title_size=18, content_size=14):
    """添加卡片式内容块"""
    # 背景
    card = add_bg_rect(slide, x, y, w, h, bg_color)
    card.shadow.inherit = False
    
    # 标题
    add_text_box(slide, x + Inches(0.2), y + Inches(0.15), w - Inches(0.4),
                 Inches(0.5), title, font_size=title_size, color=title_color,
                 bold=True, alignment=PP_ALIGN.CENTER)
    
    # 内容
    add_bullet_list(slide, x + Inches(0.2), y + Inches(0.6), w - Inches(0.4),
                    h - Inches(0.8), content_lines, font_size=content_size,
                    color=content_color, bullet_color=title_color,
                    line_spacing=1.5)

def add_section_header(slide, section_num, section_title, subtitle=""):
    """添加章节标题页"""
    # 深蓝背景
    add_bg_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(7.5), DARK_BLUE)
    
    # 红色装饰条
    add_bg_rect(slide, Inches(0), Inches(3.2), Inches(13.333), Inches(0.08), GOLD)
    
    # 章节号
    add_text_box(slide, Inches(2), Inches(2.2), Inches(2), Inches(1.2),
                 f"第{section_num}章", font_size=36, color=GOLD, bold=True,
                 alignment=PP_ALIGN.LEFT)
    
    # 标题
    add_text_box(slide, Inches(2), Inches(3.5), Inches(9), Inches(1.2),
                 section_title, font_size=44, color=WHITE, bold=True,
                 alignment=PP_ALIGN.LEFT)
    
    # 副标题
    if subtitle:
        add_text_box(slide, Inches(2), Inches(4.8), Inches(9), Inches(0.8),
                     subtitle, font_size=22, color=RGBColor(0xBB, 0xDE, 0xFB),
                     bold=False, alignment=PP_ALIGN.LEFT)

def add_content_slide(slide, title, left_items=None, right_items=None,
                      bottom_text=None, title_color=DARK_RED):
    """添加标准内容页"""
    # 顶部标题栏
    add_bg_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.1), title_color)
    add_bg_rect(slide, Inches(0), Inches(1.1), Inches(13.333), Inches(0.06), GOLD)
    
    add_text_box(slide, Inches(0.6), Inches(0.15), Inches(12), Inches(0.8),
                 title, font_size=28, color=WHITE, bold=True,
                 alignment=PP_ALIGN.LEFT)
    
    y_start = Inches(1.5)
    
    if left_items and right_items:
        # 双栏布局
        add_bg_rect(slide, Inches(0.5), y_start, Inches(5.9), Inches(5.5), LIGHT_BLUE)
        add_bg_rect(slide, Inches(6.8), y_start, Inches(6), Inches(5.5), LIGHT_BLUE)
        
        add_bullet_list(slide, Inches(0.7), y_start + Inches(0.2),
                        Inches(5.5), Inches(5), left_items, font_size=16,
                        line_spacing=1.7)
        
        add_bullet_list(slide, Inches(7.0), y_start + Inches(0.2),
                        Inches(5.6), Inches(5), right_items, font_size=16,
                        line_spacing=1.7)
    elif left_items:
        add_bullet_list(slide, Inches(0.8), y_start, Inches(11.5), Inches(5.5),
                        left_items, font_size=18, line_spacing=1.8)
    
    if bottom_text:
        add_text_box(slide, Inches(0.8), Inches(7.0), Inches(11.5), Inches(0.4),
                     bottom_text, font_size=12, color=MED_GRAY,
                     alignment=PP_ALIGN.RIGHT)


# ==================== 第1页：封面 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局

# 深蓝底色
add_bg_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(7.5), DARK_BLUE)

# 红色装饰条
add_bg_rect(slide, Inches(0), Inches(3.0), Inches(13.333), Inches(0.12), RED)
add_bg_rect(slide, Inches(0), Inches(5.2), Inches(13.333), Inches(0.06), GOLD)

# 标题
add_text_box(slide, Inches(1.5), Inches(1.5), Inches(10), Inches(1.5),
             "南湖区中小学（幼儿园）", font_size=40, color=GOLD, bold=True,
             alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1.5), Inches(2.5), Inches(10), Inches(0.8),
             "后勤管理干部工作实务培训", font_size=48, color=WHITE, bold=True,
             alignment=PP_ALIGN.CENTER)

# 副标题区域
add_text_box(slide, Inches(1.5), Inches(3.5), Inches(10), Inches(1.5),
             "校舍维修 · 工程管理 · 装备采购 · 校服采购 · 日常维护",
             font_size=24, color=RGBColor(0xBB, 0xDE, 0xFB), bold=False,
             alignment=PP_ALIGN.CENTER)

# 底部信息
add_text_box(slide, Inches(1.5), Inches(5.6), Inches(10), Inches(0.6),
             "嘉兴市南湖区教育体育局  后勤管理中心",
             font_size=20, color=WHITE, bold=False, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1.5), Inches(6.3), Inches(10), Inches(0.5),
             "2026年6月",
             font_size=18, color=RGBColor(0xBB, 0xDE, 0xFB), bold=False,
             alignment=PP_ALIGN.CENTER)


# ==================== 第2页：目录 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.1), DARK_RED)
add_bg_rect(slide, Inches(0), Inches(1.1), Inches(13.333), Inches(0.06), GOLD)
add_text_box(slide, Inches(0.6), Inches(0.15), Inches(12), Inches(0.8),
             "目  录", font_size=32, color=WHITE, bold=True, alignment=PP_ALIGN.LEFT)

toc_items = [
    ("第一章", "后勤管理工作概述与政策依据"),
    ("第二章", "校舍维修与校园日常维护管理"),
    ("第三章", "工程建设项目管理"),
    ("第四章", "教育装备采购管理"),
    ("第五章", "学生校服采购管理"),
    ("第六章", "安全生产与应急管理"),
    ("第七章", "廉洁自律与作风建设"),
]

for i, (num, title) in enumerate(toc_items):
    y = Inches(1.6) + Inches(i * 0.7)
    # 编号背景
    add_bg_rect(slide, Inches(1.5), y, Inches(1.2), Inches(0.55), BLUE)
    add_text_box(slide, Inches(1.5), y + Inches(0.05), Inches(1.2), Inches(0.5),
                 num, font_size=18, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    # 标题
    add_text_box(slide, Inches(3.0), y + Inches(0.05), Inches(8), Inches(0.5),
                 title, font_size=22, color=DARK_GRAY, bold=False)


# ==================== 第3页：第一章 章节页 ====================
add_section_header(prs.slides.add_slide(prs.slide_layouts[6]),
                   "一", "后勤管理工作概述与政策依据",
                   "明确职责定位  强化规范管理")


# ==================== 第4页：后勤管理职责定位 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_content_slide(slide, "一、后勤管理职责定位",
    left_items=[
        "后勤管理是学校工作的重要组成部分",
        "● 保障教育教学正常开展的基础支撑",
        "● 维护校园安全稳定的重要防线",
        "● 服务师生学习生活的关键环节",
        "● 提升学校办学品质的物质保障",
        "",
        "后勤管理干部的'四个角色'",
        "● 校园安全的'守护者'",
        "● 资源管理的'大管家'",
        "● 工程质量的'把关人'",
        "● 服务师生的'贴心人'",
    ],
    right_items=[
        "后勤管理工作的'六个坚持'",
        "● 坚持安全第一，预防为主",
        "● 坚持规范管理，依法依规",
        "● 坚持服务至上，保障有力",
        "● 坚持勤俭节约，精打细算",
        "● 坚持公开透明，阳光操作",
        "● 坚持持续改进，追求卓越",
        "",
        "工作目标",
        "● 设施设备完好率 ≥ 98%",
        "● 安全隐患整改率 100%",
        "● 师生满意度 ≥ 90%",
        "● 采购合规率 100%",
    ])


# ==================== 第5页：政策法规依据 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_content_slide(slide, "二、政策法规依据",
    left_items=[
        "国家层面法律法规",
        "● 《中华人民共和国教育法》",
        "● 《中华人民共和国义务教育法》",
        "● 《中华人民共和国建筑法》",
        "● 《中华人民共和国招标投标法》",
        "● 《中华人民共和国政府采购法》",
        "● 《中华人民共和国安全生产法》",
        "● 《建设工程质量管理条例》",
        "● 《学校食品安全与营养健康管理规定》",
    ],
    right_items=[
        "浙江省及南湖区相关政策",
        "● 《浙江省中小学校舍安全工程实施办法》",
        "● 《浙江省教育装备采购管理办法》",
        "● 《浙江省中小学生校服管理办法》",
        "● 《嘉兴市政府采购管理办法》",
        "● 《南湖区教育系统基建工程项目管理办法》",
        "● 《南湖区学校后勤管理工作规范》",
        "● 《南湖区中小学校舍维修专项资金管理办法》",
        "● 南湖区教育局年度后勤管理工作要点",
    ],
    bottom_text="提示：以上政策文件请结合南湖区教育局最新下发的文件汇编学习")


# ==================== 第6页：第二章 章节页 ====================
add_section_header(prs.slides.add_slide(prs.slide_layouts[6]),
                   "二", "校舍维修管理",
                   "筑牢安全底线  保障办学条件")


# ==================== 第7页：校舍维修管理概述 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_content_slide(slide, "一、校舍维修管理总体要求",
    left_items=[
        "维修分类",
        "● 日常维修：小型修缮，即时处理",
        "● 定期维修：周期性维护保养",
        "● 专项维修：针对特定问题的集中维修",
        "● 应急维修：突发情况的紧急处置",
        "● 大修改造：校舍结构性维修改造",
        "",
        "维修原则",
        "● 安全第一，质量为本",
        "● 及时响应，快速处置",
        "● 统筹安排，分类实施",
        "● 专款专用，规范管理",
    ],
    right_items=[
        "维修管理流程",
        "① 发现报修 → ② 现场勘查 → ③ 方案制定",
        "④ 预算审核 → ⑤ 审批立项 → ⑥ 组织实施",
        "⑦ 过程监管 → ⑧ 竣工验收 → ⑨ 资料归档",
        "",
        "南湖区管理要求",
        "● 建立校舍安全定期检查制度（每学期至少2次全面检查）",
        "● 建立校舍维修项目库，实行滚动管理",
        "● 维修项目须纳入年度预算，按程序报批",
        "● 维修资金实行专账管理、专款专用",
        "● 重大维修项目须报区教育局备案审批",
    ])


# ==================== 第8页：校舍安全检查制度 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_content_slide(slide, "二、校舍安全检查制度",
    left_items=[
        "检查频次要求",
        "● 日常巡查：每日一次（值班人员）",
        "● 周检查：每周一次（后勤负责人）",
        "● 月检查：每月一次（分管领导）",
        "● 学期检查：每学期开学前和期末各一次全面检查",
        "● 专项检查：汛期、台风季、冰雪天气前后",
        "",
        "检查重点内容",
        "● 建筑结构安全（墙体、梁柱、楼板）",
        "● 屋面防水及排水系统",
        "● 电气线路及配电设施",
        "● 消防设施设备",
        "● 供水、排水管道",
        "● 围墙、大门、护栏等附属设施",
    ],
    right_items=[
        "隐患分级管理",
        "● A级（一般隐患）：限期整改，校内解决",
        "● B级（较大隐患）：专项方案，上报备案",
        "● C级（重大隐患）：立即停用，上报区教育局",
        "",
        "南湖区具体要求",
        "● 建立'一校一档'校舍安全档案",
        "● 检查记录须拍照留存，形成书面报告",
        "● 发现C级隐患须24小时内上报",
        "● 建立隐患整改台账，实行销号管理",
        "● 委托有资质的第三方进行房屋安全鉴定",
        "",
        "台风汛期特别要求",
        "● 提前排查排水管网、屋顶防水",
        "● 加固户外设施、广告牌、树木",
        "● 储备应急物资（沙袋、抽水泵等）",
    ])


# ==================== 第9页：校舍维修项目管理 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_content_slide(slide, "三、校舍维修项目管理",
    left_items=[
        "项目立项审批",
        "● 维修预算5万元以下：学校自行审批",
        "● 维修预算5-20万元：报区教育局备案",
        "● 维修预算20万元以上：报区教育局审批",
        "● 涉及结构安全的维修一律上报审批",
        "",
        "施工单位选择",
        "● 建立合格供应商库",
        "● 小额维修可采用询价方式",
        "● 较大维修项目按规定招标",
        "● 严禁将项目肢解规避招标",
        "",
        "施工过程管理",
        "● 签订规范施工合同",
        "● 明确安全责任",
        "● 做好施工期间的师生安全防护",
    ],
    right_items=[
        "竣工验收要求",
        "● 组织相关人员现场验收",
        "● 对照合同和图纸逐项检查",
        "● 形成验收报告，参与人员签字",
        "● 验收不合格不得付款",
        "",
        "资金管理",
        "● 维修资金纳入年度预算",
        "● 按合同约定分期支付",
        "● 预留质量保证金（不低于5%）",
        "● 质保期满后经复验无问题再退还",
        "",
        "资料归档",
        "● 项目申报审批文件",
        "● 施工合同及附件",
        "● 施工前后对比照片",
        "● 验收报告及结算资料",
        "● 质保承诺书",
    ],
    bottom_text="注：具体金额标准以南湖区教育局最新文件为准")


# ==================== 第10页：第三章 章节页 ====================
add_section_header(prs.slides.add_slide(prs.slide_layouts[6]),
                   "三", "工程建设项目管理",
                   "规范建设程序  打造精品工程")


# ==================== 第11页：工程建设管理流程 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_content_slide(slide, "一、工程建设管理全流程",
    left_items=[
        "前期阶段",
        "● 项目建议书编制与申报",
        "● 可行性研究论证",
        "● 立项审批（发改部门）",
        "● 用地规划许可",
        "● 工程规划许可",
        "● 施工图设计及审查",
        "",
        "招标阶段",
        "● 编制招标文件",
        "● 发布招标公告",
        "● 开标、评标、定标",
        "● 签订施工合同",
        "● 办理施工许可证",
    ],
    right_items=[
        "施工阶段",
        "● 开工报告审批",
        "● 质量安全管理",
        "● 进度控制",
        "● 变更管理（严格审批）",
        "● 隐蔽工程验收",
        "● 阶段性检查",
        "",
        "竣工阶段",
        "● 竣工验收（五方责任主体）",
        "● 消防验收/备案",
        "● 环保验收",
        "● 档案移交",
        "● 竣工结算审计",
        "● 固定资产登记",
        "",
        "南湖区要求：所有新建、改扩建项目须报区教育局审核同意后方可实施"
    ])


# ==================== 第12页：工程质量安全管理 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_content_slide(slide, "二、工程质量与安全管理",
    left_items=[
        "质量管理要点",
        "● 落实五方责任主体质量终身责任制",
        "● 严格执行工程建设强制性标准",
        "● 材料进场检验制度",
        "● 隐蔽工程验收签字制度",
        "● 分部分项工程验收制度",
        "● 第三方质量检测",
        "",
        "安全管理要点",
        "● 施工单位安全资质审查",
        "● 安全教育培训",
        "● 施工现场封闭管理",
        "● 施工区域与教学区域物理隔离",
        "● 特种设备安全管理",
        "● 应急预案及演练",
    ],
    right_items=[
        "学校后勤干部职责",
        "● 协调各方关系，保障工程顺利推进",
        "● 监督施工单位文明施工",
        "● 关注施工对教学秩序的影响",
        "● 定期向区教育局报告工程进度",
        "● 妥善处理周边关系（居民投诉等）",
        "",
        "变更管理",
        "● 设计变更须按程序审批",
        "● 变更导致投资增加10%以上须重新报批",
        "● 严禁'先施工后补手续'",
        "● 所有变更须留存书面记录",
        "",
        "禁止行为",
        "✗ 违规肢解工程规避招标",
        "✗ 擅自扩大建设规模",
        "✗ 降低工程质量标准",
        "✗ 违规指定供应商或承包商",
    ])


# ==================== 第13页：第四章 章节页 ====================
add_section_header(prs.slides.add_slide(prs.slide_layouts[6]),
                   "四", "教育装备采购管理",
                   "阳光采购  规范操作")


# ==================== 第14页：装备采购管理要求 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_content_slide(slide, "一、装备采购管理基本要求",
    left_items=[
        "采购范围",
        "● 教学仪器设备（实验器材、音体美器材等）",
        "● 信息化设备（电脑、投影、交互智能平板等）",
        "● 办公家具及生活设施",
        "● 图书资料",
        "● 厨房设备及食堂用品",
        "● 安保、消防设备",
        "",
        "采购方式",
        "● 政府采购（达到限额标准）",
        "● 协议供货（纳入政采云目录）",
        "● 网上竞价",
        "● 竞争性磋商/谈判",
        "● 询价采购（小额）",
    ],
    right_items=[
        "采购流程",
        "① 需求论证 → ② 预算编制 → ③ 采购计划申报",
        "④ 采购方式审批 → ⑤ 实施采购 → ⑥ 合同签订",
        "⑦ 货物验收 → ⑧ 入库登记 → ⑨ 资金支付",
        "",
        "南湖区管理要求",
        "● 严格执行嘉兴市政府采购目录及标准",
        "● 所有采购须通过'政采云'平台操作",
        "● 装备配置须符合《浙江省中小学教育装备标准》",
        "● 建立装备采购需求论证制度",
        "● 严禁化整为零规避政府采购",
        "● 采购结果须在校内公示",
    ])


# ==================== 第15页：装备验收与资产管理 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_content_slide(slide, "二、装备验收与资产管理",
    left_items=[
        "验收管理",
        "● 成立验收小组（至少3人）",
        "● 对照合同逐项验收",
        "● 检查产品合格证、说明书",
        "● 进行功能测试和试运行",
        "● 形成验收报告，签字确认",
        "● 不合格产品坚决退换",
        "",
        "资产管理",
        "● 验收合格后及时入库登记",
        "● 建立固定资产台账",
        "● 贴附资产标签",
        "● 落实使用保管责任人",
    ],
    right_items=[
        "日常维护",
        "● 制定设备维护保养计划",
        "● 定期巡检，做好维护记录",
        "● 建立报修维修制度",
        "● 重要设备购买维保服务",
        "",
        "报废处置",
        "● 达到使用年限且无法修复的可申请报废",
        "● 报废须履行审批手续",
        "● 报废资产统一回收处理",
        "● 及时更新资产台账",
        "",
        "南湖区要求",
        "● 每年至少开展一次资产清查盘点",
        "● 资产数据与区教育局资产管理系统保持一致",
    ])


# ==================== 第16页：第五章 章节页 ====================
add_section_header(prs.slides.add_slide(prs.slide_layouts[6]),
                   "五", "学生校服采购管理",
                   "规范管理  放心着装")


# ==================== 第17页：校服采购管理要求 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_content_slide(slide, "一、校服采购管理要求",
    left_items=[
        "政策依据",
        "● 《浙江省中小学生校服管理办法》",
        "● 《教育部等四部门关于进一步加强中小学生校服管理工作的意见》",
        "● 南湖区校服管理相关规定",
        "",
        "基本原则",
        "● 自愿购买原则（不得强制）",
        "● 公开透明原则",
        "● 质量优先原则",
        "● 家长参与原则",
        "",
        "质量标准",
        "● 必须符合GB/T 31888-2015《中小学生校服》标准",
        "● 纤维含量、色牢度、甲醛含量等指标达标",
    ],
    right_items=[
        "采购程序",
        "① 成立校服选购委员会（家长代表≥80%）",
        "② 制定采购方案并公示",
        "③ 按规定方式确定供货企业",
        "④ 签订采购合同（明确质量标准）",
        "⑤ 校服样品送检（第三方检测机构）",
        "⑥ 验收合格后发放",
        "",
        "南湖区管理要求",
        "● 校服采购结果报区教育局备案",
        "● 建立校服质量追溯机制",
        "● 每批次校服须进行抽样送检",
        "● 建立家长投诉处理机制",
        "● 对困难学生实行校服费用减免",
        "● 校服选购委员会中家长代表比例不低于80%",
    ])


# ==================== 第18页：校服质量管理 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_content_slide(slide, "二、校服质量监督管理",
    left_items=[
        "质量检验制度",
        "● 供货企业须提供每批次质检报告",
        "● 学校须抽样送第三方检测机构检验",
        "● 检验费用由供货企业承担",
        "● 检验不合格不得发放给学生",
        "",
        "验收管理",
        "● 对照合同和样品逐件验收",
        "● 检查标识标签（号型、成分、洗涤说明）",
        "● 检查外观质量和缝制工艺",
        "● 建立验收台账",
    ],
    right_items=[
        "售后服务要求",
        "● 提供合理的调换期限（不少于30天）",
        "● 建立质量投诉快速响应机制",
        "● 质量问题无条件退换",
        "",
        "禁止行为",
        "✗ 强制学生购买校服",
        "✗ 指定特定品牌或供应商",
        "✗ 收受供货企业回扣或好处",
        "✗ 降低质量标准",
        "✗ 未经检验直接发放",
        "",
        "监督机制",
        "● 区教育局定期抽查",
        "● 家长委员会全程参与监督",
        "● 接受社会监督",
    ])


# ==================== 第19页：第六章 章节页 ====================
add_section_header(prs.slides.add_slide(prs.slide_layouts[6]),
                   "六", "校园日常维护管理",
                   "精细化管理  品质化服务")


# ==================== 第20页：日常维护管理 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_content_slide(slide, "一、校园日常维护管理",
    left_items=[
        "基础设施维护",
        "● 校舍建筑日常巡查与维护",
        "● 水电设施定期检查保养",
        "● 道路、场地、围墙维护",
        "● 绿化养护管理",
        "",
        "设施设备维护",
        "● 教学设备日常保养",
        "● 消防设施月度检查",
        "● 电梯等特种设备年检",
        "● 体育设施安全检查",
        "",
        "环境管理",
        "● 校园保洁（每日清扫）",
        "● 垃圾分类管理",
        "● 病媒生物防制（四害消杀）",
    ],
    right_items=[
        "管理制度建设",
        "● 建立日常巡查制度（有记录、有反馈）",
        "● 建立报修响应制度（一般问题24小时内处理）",
        "● 建立设施设备维护保养计划",
        "● 建立外包服务监管制度",
        "",
        "南湖区管理要求",
        "● 推行校园后勤管理标准化",
        "● 建立'日巡查、周检查、月通报'制度",
        "● 推进智慧校园后勤管理平台应用",
        "● 开展'最美校园'创建活动",
        "● 后勤管理纳入学校年度考核",
    ])


# ==================== 第21页：食堂与食品安全 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_content_slide(slide, "二、食堂与食品安全管理",
    left_items=[
        "食堂管理要求",
        "● 取得食品经营许可证",
        "● 落实校长陪餐制度",
        "● 实行明厨亮灶",
        "● 建立食品留样制度（48小时）",
        "● 从业人员持健康证上岗",
        "● 定期开展食品安全培训",
        "",
        "食材采购",
        "● 定点采购，索证索票",
        "● 建立进货查验台账",
        "● 严禁采购过期、变质食品",
        "● 大宗食材统一招标采购",
    ],
    right_items=[
        "食品安全管理",
        "● 落实'日管控、周排查、月调度'机制",
        "● 建立食品安全应急预案",
        "● 定期开展食品安全自查",
        "● 接受市场监管部门检查",
        "",
        "南湖区具体要求",
        "● 全面落实'互联网+明厨亮灶'",
        "● 食堂财务独立核算，坚持公益性原则",
        "● 建立家长参与食堂监督机制",
        "● 推行营养餐标准，关注学生营养健康",
        "● 食堂大宗食材纳入区级统一采购平台",
    ])


# ==================== 第22页：第七章 章节页 ====================
add_section_header(prs.slides.add_slide(prs.slide_layouts[6]),
                   "七", "安全生产与应急管理",
                   "生命至上  安全第一")


# ==================== 第23页：安全生产管理 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_content_slide(slide, "一、安全生产管理",
    left_items=[
        "安全责任体系",
        "● 落实'党政同责、一岗双责'",
        "● 校长是第一责任人",
        "● 签订安全责任书（层层落实）",
        "● 建立安全工作台账",
        "",
        "重点领域安全",
        "● 消防安全（灭火器、消防栓、疏散通道）",
        "● 用电安全（线路检查、漏电保护）",
        "● 燃气安全（食堂燃气报警装置）",
        "● 特种设备安全（电梯、锅炉）",
        "● 建筑施工安全",
    ],
    right_items=[
        "隐患排查治理",
        "● 建立隐患排查治理制度",
        "● 发现隐患立即整改",
        "● 重大隐患挂牌督办",
        "● 隐患整改闭环管理",
        "",
        "应急管理",
        "● 制定完善各类应急预案",
        "● 每学期至少开展2次应急演练",
        "● 储备应急物资",
        "● 建立应急值班制度",
        "● 突发事件及时上报（不得瞒报、迟报）",
        "",
        "南湖区要求",
        "● 落实校园安全'网格化'管理",
        "● 建立安全隐患'清单化'管理",
    ])


# ==================== 第24页：第八章 章节页 ====================
add_section_header(prs.slides.add_slide(prs.slide_layouts[6]),
                   "八", "廉洁自律与作风建设",
                   "清正廉洁  阳光后勤")


# ==================== 第25页：廉洁自律要求 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_content_slide(slide, "一、廉洁自律与作风建设",
    left_items=[
        "纪律红线",
        "✗ 严禁在采购中收受回扣、好处费",
        "✗ 严禁违规指定供应商或承包商",
        "✗ 严禁虚报冒领、截留挪用资金",
        "✗ 严禁在工程中偷工减料、以次充好",
        "✗ 严禁违规接受供应商宴请、礼品",
        "✗ 严禁泄露招标评审等保密信息",
        "",
        "廉政风险防控",
        "● 梳理岗位廉政风险点",
        "● 制定防控措施",
        "● 关键岗位定期轮岗",
    ],
    right_items=[
        "阳光操作要求",
        "● 采购信息公开发布",
        "● 采购结果校内公示",
        "● 重大事项集体决策",
        "● 接受纪检部门监督",
        "● 畅通举报投诉渠道",
        "",
        "作风建设",
        "● 强化服务意识，提升服务品质",
        "● 提高工作效率，做到快速响应",
        "● 坚持勤俭节约，反对铺张浪费",
        "● 加强业务学习，提升专业能力",
        "● 主动接受监督，保持清正廉洁",
    ])


# ==================== 第26页：总结 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(7.5), DARK_BLUE)
add_bg_rect(slide, Inches(0), Inches(3.5), Inches(13.333), Inches(0.08), GOLD)

add_text_box(slide, Inches(1.5), Inches(1.0), Inches(10), Inches(1.0),
             "总结与要求", font_size=44, color=WHITE, bold=True,
             alignment=PP_ALIGN.CENTER)

summary_items = [
    "依法依规，规范操作——所有工作必须在法律法规框架内开展",
    "安全第一，质量为本——筑牢校园安全防线，打造精品工程",
    "阳光透明，廉洁自律——让权力在阳光下运行，守住廉洁底线",
    "服务师生，保障有力——以提升师生满意度为工作目标",
    "持续学习，不断提升——紧跟政策变化，提高业务能力",
]

for i, item in enumerate(summary_items):
    y = Inches(3.9) + Inches(i * 0.65)
    add_text_box(slide, Inches(2.5), y, Inches(8.5), Inches(0.6),
                 item, font_size=22, color=RGBColor(0xBB, 0xDE, 0xFB),
                 bold=False, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1.5), Inches(7.0), Inches(10), Inches(0.4),
             "让我们共同努力，打造安全、规范、高效的校园后勤管理体系！",
             font_size=18, color=GOLD, bold=True, alignment=PP_ALIGN.CENTER)


# ==================== 保存 ====================
output_path = "/home/admin/.openclaw/workspace/南湖区中小学幼儿园后勤管理干部培训.pptx"
prs.save(output_path)
print(f"PPT已生成: {output_path}")
