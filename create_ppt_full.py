#!/usr/bin/env python3
"""
创建完整的PPT演示文稿 - 多页标准格式
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# 输出PPT路径
output_pptx = "/home/admin/.openclaw/workspace/老师您费心了_完整PPT.pptx"

# 创建PPT
prs = Presentation()
prs.slide_width = Inches(13.333)  # 16:9 比例
prs.slide_height = Inches(7.5)

# 颜色定义
primary_color = RGBColor(51, 51, 51)      # 深灰
accent_color = RGBColor(0, 122, 204)       # 蓝色
gold_color = RGBColor(255, 215, 0)         # 金色
light_bg = RGBColor(245, 245, 245)         # 浅灰背景
white = RGBColor(255, 255, 255)

def add_title_slide(prs, title, subtitle=""):
    """添加标题页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 背景
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = light_bg
    shape.line.fill.background()
    
    # 主标题
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11.333), Inches(2))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = title
    title_para.font.size = Pt(54)
    title_para.font.bold = True
    title_para.font.color.rgb = primary_color
    title_para.alignment = PP_ALIGN.CENTER
    
    # 副标题
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(11.333), Inches(1))
        sub_frame = sub_box.text_frame
        sub_para = sub_frame.paragraphs[0]
        sub_para.text = subtitle
        sub_para.font.size = Pt(28)
        sub_para.font.color.rgb = accent_color
        sub_para.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, content_list):
    """添加内容页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(1))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = title
    title_para.font.size = Pt(40)
    title_para.font.bold = True
    title_para.font.color.rgb = primary_color
    
    # 内容
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(11.933), Inches(5.5))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    
    for i, item in enumerate(content_list):
        if i == 0:
            para = content_frame.paragraphs[0]
        else:
            para = content_frame.add_paragraph()
        para.text = item
        para.font.size = Pt(24)
        para.font.color.rgb = primary_color
        para.space_after = Pt(20)
    
    return slide

def add_quote_slide(prs, quote_text, author=""):
    """添加金句页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 背景
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = primary_color
    shape.line.fill.background()
    
    # 引号
    quote_box = slide.shapes.add_textbox(Inches(1.5), Inches(2), Inches(10.333), Inches(3))
    quote_frame = quote_box.text_frame
    quote_para = quote_frame.paragraphs[0]
    quote_para.text = f'"{quote_text}"'
    quote_para.font.size = Pt(44)
    quote_para.font.color.rgb = gold_color
    quote_para.alignment = PP_ALIGN.CENTER
    quote_frame.word_wrap = True
    
    # 作者
    if author:
        author_box = slide.shapes.add_textbox(Inches(1.5), Inches(5), Inches(10.333), Inches(1))
        author_frame = author_box.text_frame
        author_para = author_frame.paragraphs[0]
        author_para.text = f"— {author}"
        author_para.font.size = Pt(28)
        author_para.font.color.rgb = white
        author_para.alignment = PP_ALIGN.CENTER
    
    return slide

# ============ 创建幻灯片 ============

# 第1页：封面
add_title_slide(prs, "老师您费心了", "五种父母心态类型解析")

# 第2页：引言
add_content_slide(prs, "为什么是「老师您费心了」？", [
    "• 这句话是家校沟通中最常见的客套话",
    "• 但背后隐藏着不同的父母心态",
    "• 理解这些心态，才能真正走进家长内心",
    "• 今天不为「说服」，只为「交心」"
])

# 第3页：类型一
slide = add_content_slide(prs, "【释然型】父母", [
    "心态特征：",
    "• 接受孩子的普通",
    "• 平安健康就好",
    "• 不再过度焦虑成绩",
    "",
    "典型话语：",
    "• 「孩子已经尽力了」",
    "• 「平安健康就好」",
    "",
    "沟通建议：",
    "• 肯定家长的心态转变",
    "• 分享孩子其他方面的闪光点"
])

# 第4页：类型二
slide = add_content_slide(prs, "【心累型】父母", [
    "心态特征：",
    "• 管不动了，选择放手",
    "• 自己也感到疲惫",
    "• 需要自我和解",
    "",
    "典型话语：",
    "• 「管不动了，随他去吧」",
    "• 「我也需要放过自己」",
    "",
    "沟通建议：",
    "• 表达理解与共情",
    "• 提供具体可操作的小建议"
])

# 第5页：类型三
slide = add_content_slide(prs, "【保护型】父母", [
    "心态特征：",
    "• 心疼孩子压力大",
    "• 不想再施压",
    "• 站在孩子这边",
    "",
    "典型话语：",
    "• 「孩子不容易」",
    "• 「不想再给他增加压力」",
    "",
    "沟通建议：",
    "• 肯定家长的保护欲",
    "• 说明适度压力的必要性"
])

# 第6页：类型四
slide = add_content_slide(prs, "【防备型】父母", [
    "心态特征：",
    "• 对老师有戒备心理",
    "• 不想被频繁打扰",
    "• 可能有过负面经历",
    "",
    "典型话语：",
    "• 「老师别总找我了」",
    "• 「我们自己也头疼」",
    "",
    "沟通建议：",
    "• 建立信任是关键",
    "• 先报喜再报忧"
])

# 第7页：类型五
slide = add_content_slide(prs, "【客套型】父母", [
    "心态特征：",
    "• 表面客气，内心不认同",
    "• 嘴上说费心，回家该骂还骂",
    "• 缺乏真正的配合",
    "",
    "典型话语：",
    "• 「老师您费心了」",
    "• 「回家我们会教育的」",
    "",
    "沟通建议：",
    "• 用具体案例打动家长",
    "• 邀请参与学校活动"
])

# 第8页：对比总结
add_content_slide(prs, "五种类型对比", [
    "释然型 → 接纳与放下",
    "心累型 → 疲惫与无奈",
    "保护型 → 心疼与守护",
    "防备型 → 戒备与疏离",
    "客套型 → 表面与敷衍",
    "",
    "理解是沟通的第一步"
])

# 第9页：金句
add_quote_slide(prs, "今天我不为「说服」，只为交心", "教育的本质是理解")

# 第10页：结束页
add_title_slide(prs, "谢谢聆听", "期待与您真诚沟通")

# 保存PPT
prs.save(output_pptx)
print(f"完整PPT已生成：{output_pptx}")
print(f"共 10 页幻灯片")
