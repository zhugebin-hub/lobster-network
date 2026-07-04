#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
关键帧应用教学案例 PPT 生成器
"""

from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Cm(33.867)
prs.slide_height = Cm(19.05)

COLORS = {
    'primary': (41, 128, 185),
    'secondary': (52, 152, 219),
    'accent': (230, 126, 34),
    'dark': (44, 62, 80),
    'light': (236, 240, 241),
    'white': (255, 255, 255),
}

def set_background(slide, color=COLORS['white']):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*color)

def add_header(slide, title):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(0), Cm(0), prs.slide_width, Cm(1.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*COLORS['primary'])
    shape.line.fill.background()
    
    textbox = slide.shapes.add_textbox(Cm(1.5), Cm(0.3), Cm(25), Cm(1))
    tf = textbox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = RGBColor(*COLORS['white'])
    p.font.name = 'Microsoft YaHei'

def add_footer(slide, page_num=""):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(2), prs.slide_height - Cm(1), prs.slide_width - Cm(4), Cm(0.05))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*COLORS['secondary'])
    shape.line.fill.background()
    
    textbox = slide.shapes.add_textbox(prs.slide_width - Cm(3), prs.slide_height - Cm(1.3), Cm(2), Cm(0.5))
    tf = textbox.text_frame
    p = tf.paragraphs[0]
    p.text = page_num
    p.font.size = Pt(11)
    p.font.color.rgb = RGBColor(*COLORS['dark'])
    p.font.name = 'Microsoft YaHei'
    p.alignment = PP_ALIGN.RIGHT

def add_title_slide(prs, title, subtitle, footer=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(0), Cm(0), prs.slide_width, Cm(5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*COLORS['primary'])
    shape.line.fill.background()
    
    textbox = slide.shapes.add_textbox(Cm(2), Cm(1.5), Cm(28), Cm(2.5))
    tf = textbox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(*COLORS['white'])
    p.font.name = 'Microsoft YaHei'
    
    textbox = slide.shapes.add_textbox(Cm(2), Cm(4.5), Cm(20), Cm(1.5))
    tf = textbox.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(22)
    p.font.color.rgb = RGBColor(*COLORS['white'])
    p.font.name = 'Microsoft YaHei'
    
    if footer:
        textbox = slide.shapes.add_textbox(Cm(2), prs.slide_height - Cm(3), Cm(15), Cm(1))
        tf = textbox.text_frame
        p = tf.paragraphs[0]
        p.text = footer
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(*COLORS['light'])
        p.font.name = 'Microsoft YaHei'
    
    return slide

def add_content_slide(prs, title, bullets, page_num=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_header(slide, title)
    
    y = Cm(2.5)
    for bullet in bullets:
        textbox = slide.shapes.add_textbox(Cm(2.5), y, Cm(29), Cm(0.7))
        tf = textbox.text_frame
        p = tf.paragraphs[0]
        p.text = bullet
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(*COLORS['dark'])
        p.font.name = 'Microsoft YaHei'
        y += Cm(0.9)
    
    add_footer(slide, page_num)
    return slide

def add_diagram_slide(prs, title, items, page_num=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_header(slide, title)
    
    y = Cm(3)
    for i, (item_title, item_desc) in enumerate(items):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(2), y, Cm(30), Cm(2))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*COLORS['light'])
        shape.line.color.rgb = RGBColor(*COLORS['primary'])
        shape.line.width = Pt(2)
        
        textbox = slide.shapes.add_textbox(Cm(2.5), y + Cm(0.3), Cm(27), Cm(0.6))
        tf = textbox.text_frame
        p = tf.paragraphs[0]
        p.text = item_title
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = RGBColor(*COLORS['primary'])
        p.font.name = 'Microsoft YaHei'
        
        textbox = slide.shapes.add_textbox(Cm(2.5), y + Cm(1), Cm(27), Cm(0.8))
        tf = textbox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = item_desc
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(*COLORS['dark'])
        p.font.name = 'Microsoft YaHei'
        
        y += Cm(2.3)
    
    add_footer(slide, page_num)
    return slide

def add_comparison_slide(prs, title, headers, rows, page_num=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_header(slide, title)
    
    col_width = Cm(10)
    row_height = Cm(1.5)
    start_x = Cm(2)
    start_y = Cm(3)
    
    for i, header in enumerate(headers):
        x = start_x + i * col_width
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, start_y, col_width, Cm(1.2))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*COLORS['primary'])
        shape.line.fill.background()
        
        textbox = slide.shapes.add_textbox(x + Cm(0.3), start_y + Cm(0.2), col_width - Cm(0.6), Cm(0.8))
        tf = textbox.text_frame
        p = tf.paragraphs[0]
        p.text = header
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = RGBColor(*COLORS['white'])
        p.font.name = 'Microsoft YaHei'
        p.alignment = PP_ALIGN.CENTER
    
    for row_idx, row_data in enumerate(rows):
        y = start_y + Cm(1.2) + row_idx * row_height
        for col_idx, cell in enumerate(row_data):
            x = start_x + col_idx * col_width
            shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, col_width - Cm(0.1), row_height - Cm(0.1))
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(*COLORS['light'])
            shape.line.color.rgb = RGBColor(*COLORS['dark'])
            shape.line.width = Pt(1)
            
            textbox = slide.shapes.add_textbox(x + Cm(0.3), y + Cm(0.3), col_width - Cm(0.9), row_height - Cm(0.6))
            tf = textbox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = cell
            p.font.size = Pt(13)
            p.font.color.rgb = RGBColor(*COLORS['dark'])
            p.font.name = 'Microsoft YaHei'
    
    add_footer(slide, page_num)
    return slide

def add_end_slide(prs, title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(0), Cm(0), prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*COLORS['primary'])
    shape.line.fill.background()
    
    textbox = slide.shapes.add_textbox(Cm(6), Cm(7), prs.slide_width - Cm(12), Cm(3))
    tf = textbox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = RGBColor(*COLORS['white'])
    p.font.name = 'Microsoft YaHei'
    p.alignment = PP_ALIGN.CENTER
    
    if subtitle:
        textbox = slide.shapes.add_textbox(Cm(6), Cm(10), prs.slide_width - Cm(12), Cm(2))
        tf = textbox.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(22)
        p.font.color.rgb = RGBColor(*COLORS['white'])
        p.font.name = 'Microsoft YaHei'
        p.alignment = PP_ALIGN.CENTER
    
    return slide

# 生成幻灯片
add_title_slide(prs, "关键帧技术的应用", "数字视音频制作课程", "绍兴柯桥职业学校 · 金国平")

add_content_slide(prs, "教学目标", [
    "【知识目标】",
    "  • 理解关键帧的概念和作用",
    "  • 掌握关键帧的创建方法",
    "  • 了解关键帧插值类型（线性/贝塞尔/定格）",
    "",
    "【能力目标】",
    "  • 能够独立创建和编辑关键帧",
    "  • 能够制作简单的位移动画",
    "  • 能够调节关键帧曲线控制动画节奏",
    "",
    "【素养目标】",
    "  • 培养精益求精的工匠精神",
    "  • 提升审美能力和创意思维",
], "01/12")

add_content_slide(prs, "教学重难点", [
    "【教学重点】",
    "  • 关键帧的创建与编辑",
    "  • 关键帧插值类型应用",
    "",
    "【教学难点】",
    "  • 贝塞尔曲线调节",
    "  • 动画节奏把握",
    "",
    "【突破策略】",
    "  • 演示 + 实操 + 个别指导",
    "  • 对比演示 + 案例分析",
    "  • 分步演示 + 练习巩固",
], "02/12")

add_content_slide(prs, "案例：校园宣传片头动画", [
    "【案例描述】",
    "制作一个 10 秒的校园宣传片头，包含：",
    "",
    "  • 校名文字从左侧飞入（0-2 秒）",
    "  • 校徽从放大到正常（2-4 秒）",
    "  • 标语文字淡入显示（4-6 秒）",
    "  • 背景颜色渐变变化（6-8 秒）",
    "  • 动画完成，定格 2 秒（8-10 秒）",
], "03/12")

add_diagram_slide(prs, "教学过程（90 分钟）", [
    ("情境导入（10 分钟）", "播放优秀作品，引出关键帧技术"),
    ("知识讲解（20 分钟）", "关键帧概念、插值类型、速度曲线"),
    ("案例实操（40 分钟）", "三个任务：校名飞入/校徽缩放/标语淡入"),
    ("作品展示与评价（15 分钟）", "展示作品，多元评价"),
    ("课堂小结（5 分钟）", "知识梳理，布置作业"),
], "04/12")

add_content_slide(prs, "知识点 1：什么是关键帧", [
    "【定义】",
    "动画中对象属性发生变化的帧",
    "",
    "【工作原理】",
    "软件自动计算关键帧之间的过渡（补间动画）",
    "",
    "【生活类比】",
    "翻书动画——关键页 + 自动补间",
    "",
    "【演示操作】",
    "1. 新建合成（1920×1080，10 秒）",
    "2. 创建文字图层",
    "3. 设置位置关键帧",
    "4. 播放预览",
], "05/12")

add_comparison_slide(prs, "知识点 2：关键帧插值类型",
    ["插值类型", "特点", "适用场景"],
    [
        ["线性插值", "匀速运动，机械感", "机械运动、匀速移动"],
        ["贝塞尔插值", "缓入缓出，自然流畅", "大多数自然运动"],
        ["定格插值", "瞬间变化，无过渡", "闪烁、切换效果"],
    ],
    "06/12"
)

add_content_slide(prs, "知识点 3：速度曲线调节", [
    "【打开速度曲线】",
    "选择关键帧，按 F9 转换为贝塞尔插值",
    "",
    "【调节方法】",
    "  • 打开曲线编辑器",
    "  • 调节入点和出点手柄",
    "  • 观察动画节奏变化",
    "",
    "【缓入缓出原理】",
    "开始慢→中间快→结束慢，符合自然运动规律",
], "07/12")

add_content_slide(prs, "任务一：校名飞入动画（15 分钟）", [
    "【操作步骤】",
    "1. 新建文字图层，输入学校名称",
    "2. 第 0 帧：位置（-500, 540），不透明度 0%",
    "3. 第 1 秒：不透明度 100%",
    "4. 第 2 秒：位置（960, 540）",
    "5. 选择所有关键帧，按 F9 添加缓动",
    "6. 调节速度曲线，使飞入更自然",
    "",
    "【技术要点】",
    "  • 位置关键帧控制位移",
    "  • 不透明度关键帧控制淡入",
], "08/12")

add_content_slide(prs, "任务二：校徽缩放动画（15 分钟）", [
    "【操作步骤】",
    "1. 导入校徽素材（PNG 格式）",
    "2. 第 2 秒：缩放（0%），不透明度 0%",
    "3. 第 3 秒：不透明度 100%",
    "4. 第 4 秒：缩放（100%）",
    "5. 添加弹性效果（可选）",
    "",
    "【技术要点】",
    "  • 锚点位置影响缩放中心",
    "  • 缩放动画注意比例锁定",
    "  • 弹性效果增加动画趣味",
], "09/12")

add_content_slide(prs, "任务三：标语淡入动画（10 分钟）", [
    "【操作步骤】",
    "1. 创建文字图层，输入宣传标语",
    "2. 第 4 秒：不透明度 0%",
    "3. 第 6 秒：不透明度 100%",
    "4. 添加模糊效果（可选）",
    "5. 调节动画节奏",
    "",
    "【技术要点】",
    "  • 透明度关键帧控制淡入淡出",
    "  • 模糊效果增加层次感",
], "10/12")

add_content_slide(prs, "作品评价标准", [
    "【评价维度】",
    "  • 技术运用（40%）：关键帧使用正确，曲线调节合理",
    "  • 动画效果（30%）：动画流畅，节奏舒适",
    "  • 创意设计（20%）：有创意，视觉效果 good",
    "  • 规范操作（10%）：图层命名规范，文件管理有序",
    "",
    "【评价方式】",
    "  • 学生自评：反思学习过程",
    "  • 同伴互评：学习他人优点",
    "  • 教师点评：总结共性问题",
], "11/12")

add_end_slide(prs, "感谢观看", "敬请批评指正 · Q&A")

output_path = '/home/admin/.openclaw/workspace/关键帧应用教学 PPT.pptx'
prs.save(output_path)
print(f"PPT 已生成：{output_path}")
print(f"共 {len(prs.slides)} 页")
