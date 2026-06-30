#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Generate Full PPT for "AI Agents as Digital Employees: Reconstructing Smart Home Scenarios"
30 slides complete version
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color scheme
PRIMARY = RGBColor(0x1A, 0x73, 0xE8)
DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)
LIGHT_BG = RGBColor(0xF8, 0xF9, 0xFA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x33, 0x33, 0x33)
GRAY = RGBColor(0x66, 0x66, 0x66)
ACCENT = RGBColor(0x00, 0xBC, 0xD4)
SUCCESS = RGBColor(0x00, 0xC8, 0x53)
WARNING = RGBColor(0xFF, 0x91, 0x00)
ORANGE = RGBColor(0xFF, 0x6D, 0x00)
LIGHT_GRAY = RGBColor(0xF0, 0xF0, 0xF0)

def set_slide_bg(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape_bg(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18, 
                 color=BLACK, bold=False, alignment=PP_ALIGN.LEFT, 
                 font_name='微软雅黑', line_spacing=1.2):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.line_spacing = Pt(font_size * line_spacing)
    return txBox

def add_bullet_list(slide, left, top, width, height, items, font_size=16, 
                    color=BLACK, font_name='微软雅黑', spacing=1.5, bullet_char='\u2022'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"{bullet_char} {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.line_spacing = Pt(font_size * spacing)
        p.space_after = Pt(4)
    return txBox

def add_title_bar(slide, title_text, subtitle_text=None):
    bar = add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.3), PRIMARY)
    add_text_box(slide, Inches(0.8), Inches(0.25), Inches(11.5), Inches(0.8),
                 title_text, font_size=32, color=WHITE, bold=True, alignment=PP_ALIGN.LEFT)
    if subtitle_text:
        add_text_box(slide, Inches(0.8), Inches(0.85), Inches(11.5), Inches(0.4),
                     subtitle_text, font_size=16, color=RGBColor(0xE0, 0xE0, 0xE0), 
                     alignment=PP_ALIGN.LEFT)

def add_slide_number(slide, num, total=30):
    add_text_box(slide, Inches(12.3), Inches(7.0), Inches(0.8), Inches(0.4),
                 f"{num}/{total}", font_size=12, color=GRAY, alignment=PP_ALIGN.RIGHT)

def add_table(slide, left, top, width, height, headers, rows, header_color=PRIMARY):
    cols = len(headers)
    table_shape = slide.shapes.add_table(len(rows) + 1, cols, left, top, width, height)
    table = table_shape.table
    
    for col in range(cols):
        table.columns[col].width = width // cols
    
    for j, header in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = header
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(13)
            paragraph.font.bold = True
            paragraph.font.color.rgb = WHITE
            paragraph.font.name = '微软雅黑'
            paragraph.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_color
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    for i, row in enumerate(rows):
        for j, value in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = str(value)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(12)
                paragraph.font.color.rgb = BLACK
                paragraph.font.name = '微软雅黑'
                paragraph.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if i % 2 == 0 else RGBColor(0xF5, 0xF5, 0xF5)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

# ========== Generate Slides ==========

# Slide 1: Cover
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_shape_bg(slide, Inches(0), Inches(3.2), Inches(13.333), Inches(0.06), ACCENT)
add_shape_bg(slide, Inches(0), Inches(4.8), Inches(13.333), Inches(0.03), PRIMARY)
add_text_box(slide, Inches(1.5), Inches(1.5), Inches(10.3), Inches(1.5),
             "智能体作为数字员工", font_size=48, color=WHITE, bold=True,
             alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.5), Inches(3.5), Inches(10.3), Inches(1.0),
             "重构智能家居场景", font_size=36, color=ACCENT, bold=True,
             alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.5), Inches(5.2), Inches(10.3), Inches(0.6),
             "从人控制设备到设备自主协作", font_size=20, color=GRAY,
             alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.5), Inches(6.3), Inches(10.3), Inches(0.4),
             "初识人工智能  |  诸葛斌 / 虾尔 AI  |  2026年5月",
             font_size=14, color=GRAY, alignment=PP_ALIGN.CENTER)
add_slide_number(slide, 1)

# Slide 2: Table of Contents
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, LIGHT_BG)
add_title_bar(slide, "目录", "CONTENTS")

toc_items = [
    ("01", "什么是数字员工？", "核心概念与演进阶段"),
    ("02", "场景一：下班回家", "完整智能体协作链"),
    ("03", "场景二：设备全生命周期", "购买到报废的全流程"),
    ("04", "技术架构", "系统架构与关键技术"),
    ("05", "挑战与思考", "技术、产业、伦理挑战"),
    ("06", "产业与教育", "国家战略与高校机会"),
]

y = Inches(1.8)
for num, title, desc in toc_items:
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.2), y, Inches(0.6), Inches(0.6))
    circle.fill.solid()
    circle.fill.fore_color.rgb = PRIMARY
    circle.line.fill.background()
    tf = circle.text_frame
    tf.paragraphs[0].text = num
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    add_text_box(slide, Inches(2.1), y + Inches(0.05), Inches(4), Inches(0.35),
                 title, font_size=20, color=BLACK, bold=True)
    add_text_box(slide, Inches(2.1), y + Inches(0.38), Inches(8), Inches(0.25),
                 desc, font_size=14, color=GRAY)
    y += Inches(0.8)
add_slide_number(slide, 2)

# Slide 3: Core Concept - Digital Employee
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, LIGHT_BG)
add_title_bar(slide, "核心概念：数字员工", "不是工具，是同事")

headers = ["维度", "传统模式", "数字员工模式"]
rows = [
    ["交互方式", "点击/滑动/表单", "自然语言对话"],
    ["控制主体", "人主动操作", "智能体自主执行"],
    ["设备关系", "孤岛式App", "群聊即OS"],
    ["身份模型", "匿名设备", "有身份、可追溯"],
]
add_table(slide, Inches(1.0), Inches(1.8), Inches(11.3), Inches(3.2), headers, rows)

add_shape_bg(slide, Inches(1.0), Inches(5.5), Inches(11.3), Inches(1.2), PRIMARY)
add_text_box(slide, Inches(1.5), Inches(5.7), Inches(10.3), Inches(0.8),
             "一切皆是对话 —— 对话是人类最低摩擦的交互方式。",
             font_size=22, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_slide_number(slide, 3)

# Slide 4: Evolution Stages
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, LIGHT_BG)
add_title_bar(slide, "智能终端生态演进", "设备获得人格的三个阶段")

stages = [
    ("1", "联网化", "已完成", "设备能连WiFi/蓝牙\n但需要专属App", PRIMARY),
    ("2", "智能体化", "进行中", "设备拥有独立数字身份\n成为群聊成员", ACCENT),
    ("3", "自主化", "未来", "设备之间也能对话协作\n无需人类介入", ORANGE),
]

x = Inches(0.8)
for num, title, status, desc, color in stages:
    card = add_shape_bg(slide, x, Inches(2.0), Inches(3.6), Inches(4.5), WHITE)
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(1.3), Inches(2.2), Inches(1.0), Inches(1.0))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    tf = circle.text_frame
    tf.paragraphs[0].text = num
    tf.paragraphs[0].font.size = Pt(28)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    add_text_box(slide, x + Inches(0.3), Inches(3.4), Inches(3.0), Inches(0.5),
                 title, font_size=24, color=BLACK, bold=True, alignment=PP_ALIGN.CENTER)
    
    status_shape = add_shape_bg(slide, x + Inches(0.8), Inches(4.0), Inches(2.0), Inches(0.35), color)
    tf = status_shape.text_frame
    tf.paragraphs[0].text = status
    tf.paragraphs[0].font.size = Pt(12)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    add_text_box(slide, x + Inches(0.3), Inches(4.6), Inches(3.0), Inches(1.5),
                 desc, font_size=16, color=GRAY, alignment=PP_ALIGN.CENTER, line_spacing=1.6)
    
    if x < Inches(8.5):
        add_text_box(slide, x + Inches(3.6), Inches(3.8), Inches(0.8), Inches(0.5),
                     "\u2192", font_size=28, color=GRAY, alignment=PP_ALIGN.CENTER)
    x += Inches(4.0)
add_slide_number(slide, 4)

# Slide 5: Scenario 1 - Coming Home
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, LIGHT_BG)
add_title_bar(slide, "场景一：下班回家", "完整智能体协作链")

info_items = [
    "时间：周五 18:30",
    "地点：杭州某智能社区",
    "群聊：我家钉钉群",
    "成员：主人 + 12个家庭智能体",
]
add_bullet_list(slide, Inches(1.0), Inches(1.8), Inches(11.0), Inches(2.0),
                info_items, font_size=20, color=BLACK, spacing=1.6)

add_shape_bg(slide, Inches(2.0), Inches(4.5), Inches(9.3), Inches(2.0), PRIMARY)
add_text_box(slide, Inches(2.5), Inches(4.8), Inches(8.3), Inches(0.6),
             "你只说了一句：", font_size=18, color=RGBColor(0xE0, 0xE0, 0xE0),
             alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(2.5), Inches(5.3), Inches(8.3), Inches(1.0),
             "下班了，回家。", font_size=36, color=WHITE, bold=True,
             alignment=PP_ALIGN.CENTER)
add_slide_number(slide, 5)

# Slide 6: Act 1 - Intent
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, LIGHT_BG)
add_title_bar(slide, "第一幕：意图发起", "18:30")

add_shape_bg(slide, Inches(2.0), Inches(2.0), Inches(9.3), Inches(4.5), WHITE)
add_text_box(slide, Inches(2.5), Inches(2.5), Inches(8.3), Inches(0.6),
             "你发消息：", font_size=20, color=GRAY, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(2.5), Inches(3.2), Inches(8.3), Inches(1.0),
             "下班了，回家。", font_size=36, color=PRIMARY, bold=True,
             alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(2.5), Inches(4.5), Inches(8.3), Inches(1.5),
             "没有指定温度\n没有说几点到\n没有要求任何设备操作",
             font_size=18, color=GRAY, alignment=PP_ALIGN.CENTER, line_spacing=1.8)

add_slide_number(slide, 6)

# Slide 7: Act 2 - Travel Response
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, LIGHT_BG)
add_title_bar(slide, "第二幕：出行智能体响应", "18:30:02 - 18:30:05")

# Car
add_shape_bg(slide, Inches(1.0), Inches(1.8), Inches(5.0), Inches(2.0), WHITE)
add_text_box(slide, Inches(1.3), Inches(1.9), Inches(4.4), Inches(0.4),
             "🚗 汽车智能体", font_size=18, color=PRIMARY, bold=True)
add_text_box(slide, Inches(1.3), Inches(2.4), Inches(4.4), Inches(1.2),
             "已从地下车库B2-045出发\n至电梯口待命\n已到达电梯口。预计19:12到家。",
             font_size=15, color=BLACK, line_spacing=1.6)

# Navigation
add_shape_bg(slide, Inches(6.5), Inches(1.8), Inches(5.8), Inches(2.0), WHITE)
add_text_box(slide, Inches(6.8), Inches(1.9), Inches(5.2), Inches(0.4),
             "🗺️ 导航智能体", font_size=18, color=ACCENT, bold=True)
add_text_box(slide, Inches(6.8), Inches(2.4), Inches(5.2), Inches(1.2),
             "当前路线42分钟（德胜高架拥堵）\n备选路线38分钟（秋石高架）\n建议走秋石，已规划最优方案。",
             font_size=15, color=BLACK, line_spacing=1.6)

add_text_box(slide, Inches(1.0), Inches(4.5), Inches(11.3), Inches(0.4),
             "技术细节：自动唤醒车辆 → 调用导航API → 广播ETA到群聊",
             font_size=16, color=GRAY, alignment=PP_ALIGN.CENTER)
add_slide_number(slide, 7)

# Slide 8: Act 3 - Home Negotiation
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, LIGHT_BG)
add_title_bar(slide, "第三幕：家居智能体自主协商", "关键转变：从人控制到设备自己商量")

add_text_box(slide, Inches(1.0), Inches(1.8), Inches(11.3), Inches(0.6),
             "🏠 家庭中枢：主人预计19:12到家。当前室外35°C，室内31°C。启动归家准备流程。",
             font_size=18, color=PRIMARY, bold=True)

# Device cards
devices = [
    ("❄️ 空调", "31°C→26°C\n18:55启动最优\n电费0.8元", PRIMARY),
    ("🧊 冰箱", "三文鱼解冻\n19:00前移至冷藏\n通知厨电准备", ACCENT),
    ("🚪 门窗", "19:00关窗\n19:12解锁\n安全提醒", SUCCESS),
    ("💡 灯光", "18:45玄关灯\n19:10暖光\n引导路径", ORANGE),
    ("🌡️ 热水器", "根据历史数据\n20:00启动加热\n避免浪费", PRIMARY),
]

x = Inches(0.5)
for icon, desc, color in devices:
    card = add_shape_bg(slide, x, Inches(2.8), Inches(2.3), Inches(3.5), WHITE)
    add_shape_bg(slide, x, Inches(2.8), Inches(2.3), Inches(0.06), color)
    add_text_box(slide, x + Inches(0.1), Inches(3.0), Inches(2.1), Inches(0.4),
                 icon, font_size=18, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.1), Inches(3.5), Inches(2.1), Inches(2.5),
                 desc, font_size=14, color=BLACK, alignment=PP_ALIGN.CENTER, line_spacing=1.6)
    x += Inches(2.5)
add_slide_number(slide, 8)

# Slide 9: AC Optimization
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, LIGHT_BG)
add_title_bar(slide, "空调智能体的优化策略", "不是现在开，而是计算最优启动时间")

add_shape_bg(slide, Inches(1.0), Inches(1.8), Inches(11.3), Inches(4.5), WHITE)

add_text_box(slide, Inches(1.5), Inches(2.0), Inches(10.3), Inches(0.4),
             "❄️ 空调智能体分析：", font_size=22, color=PRIMARY, bold=True)

lines = [
    "当前室内31°C，目标26°C，需要降温5度，预计耗时25分钟",
    "",
    "优化策略：18:55启动（到家前17分钟）",
    "  → 19:12达到26°C，避免过早运行浪费电",
    "",
    "电价考虑：19:00后进入平电0.538元/度",
    "  → 18:55启动刚好利用平电时段",
    "",
    "预计电费：0.8元（30分钟运行）",
]
add_text_box(slide, Inches(1.5), Inches(2.6), Inches(10.3), Inches(3.5),
             "\n".join(lines), font_size=16, color=BLACK, line_spacing=1.5)
add_slide_number(slide, 9)

# Slide 10: Fridge Decision
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, LIGHT_BG)
add_title_bar(slide, "冰箱智能体的决策", "冰箱知道晚餐需求，主动安排解冻")

add_shape_bg(slide, Inches(1.0), Inches(1.8), Inches(11.3), Inches(4.5), WHITE)

add_text_box(slide, Inches(1.5), Inches(2.0), Inches(10.3), Inches(0.4),
             "🧊 冰箱智能体分析：", font_size=22, color=ACCENT, bold=True)

lines = [
    "库存分析：",
    "  - 冷冻区：三文鱼（剩余2天）、鸡胸肉（剩余5天）",
    "  - 冷藏区：沙拉、酸奶、鸡蛋",
    "",
    "解冻建议：三文鱼需在19:00前移至冷藏室解冻，19:30可食用",
    "",
    "执行计划：",
    "  1. 18:45 启动三文鱼解冻程序（冷藏室局部升温至2°C）",
    "  2. 通知厨电智能体：19:15准备烤箱预热",
]
add_text_box(slide, Inches(1.5), Inches(2.6), Inches(10.3), Inches(3.5),
             "\n".join(lines), font_size=16, color=BLACK, line_spacing=1.5)
add_slide_number(slide, 10)

# Slide 11: Coordination
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, LIGHT_BG)
add_title_bar(slide, "门窗/灯光/热水器协同", "多设备联动，无缝配合")

# Door/Window
add_shape_bg(slide, Inches(1.0), Inches(1.8), Inches(3.5), Inches(4.5), WHITE)
add_text_box(slide, Inches(1.2), Inches(1.9), Inches(3.1), Inches(0.4),
             "🚪 门窗智能体", font_size=18, color=SUCCESS, bold=True)
add_text_box(slide, Inches(1.2), Inches(2.5), Inches(3.1), Inches(3.5),
             "19:00 关窗（空调启动前）\n19:10 关窗帘\n19:12 大门解锁\n安全提醒：阵雨概率30%",
             font_size=15, color=BLACK, line_spacing=1.6)

# Lights
add_shape_bg(slide, Inches(4.9), Inches(1.8), Inches(3.5), Inches(4.5), WHITE)
add_text_box(slide, Inches(5.1), Inches(1.9), Inches(3.1), Inches(0.4),
             "💡 灯光智能体", font_size=18, color=ORANGE, bold=True)
add_text_box(slide, Inches(5.1), Inches(2.5), Inches(3.1), Inches(3.5),
             "18:45 玄关灯开启\n19:10 暖光氛围灯\n19:12 走廊灯渐亮\n卧室保持暗光",
             font_size=15, color=BLACK, line_spacing=1.6)

# Water Heater
add_shape_bg(slide, Inches(8.8), Inches(1.8), Inches(3.5), Inches(4.5), WHITE)
add_text_box(slide, Inches(9.0), Inches(1.9), Inches(3.1), Inches(0.4),
             "🌡️ 热水器智能体", font_size=18, color=PRIMARY, bold=True)
add_text_box(slide, Inches(9.0), Inches(2.5), Inches(3.1), Inches(3.5),
             "当前水温45°C，目标55°C\n根据历史数据，20:15洗澡\n20:00启动加热\n避免提前加热浪费电",
             font_size=15, color=BLACK, line_spacing=1.6)
add_slide_number(slide, 11)

# Slide 12: Summary Confirmation
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, LIGHT_BG)
add_title_bar(slide, "家庭中枢汇总确认", "18:31:00")

headers = ["时间", "动作", "设备", "耗电", "费用"]
rows = [
    ["18:45", "三文鱼解冻", "冰箱", "0.1度", "0.05元"],
    ["18:45", "玄关灯开启", "灯光", "0.05度", "0.03元"],
    ["18:55", "主卧关窗", "门窗", "-", "-"],
    ["19:00", "关窗+窗帘", "门窗", "-", "-"],
    ["19:12", "门锁解锁", "门窗", "-", "-"],
    ["19:15", "烤箱预热", "厨电", "0.4度", "0.22元"],
    ["19:15", "电饭煲启动", "厨电", "0.5度", "0.27元"],
    ["18:55-19:25", "空调运行", "空调", "1.0度", "0.54元"],
    ["20:00-20:10", "热水加热", "热水器", "1.5度", "0.81元"],
]
add_table(slide, Inches(1.0), Inches(1.8), Inches(11.3), Inches(4.0), headers, rows)

add_shape_bg(slide, Inches(1.0), Inches(6.0), Inches(11.3), Inches(0.8), PRIMARY)
add_text_box(slide, Inches(1.5), Inches(6.1), Inches(10.3), Inches(0.6),
             "预计总费用：约1.92元 | 总耗电：约3.55度 | 30秒后自动确认",
             font_size=18, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_slide_number(slide, 12)

# Slide 13: Dynamic Adjustment
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, LIGHT_BG)
add_title_bar(slide, "第五幕：动态调整", "实时响应变化")

# Timeline
events = [
    ("19:08", "🗺️ 导航智能体", "秋石高架事故，延迟8分钟。ETA更新：19:20", WARNING),
    ("19:08:05", "🏠 家庭中枢", "空调启动推迟至19:03，节省0.3元\n烤箱预热推迟至19:23", PRIMARY),
    ("19:18", "🚗 汽车智能体", "已到达小区入口，识别到人脸，道闸已自动开启", ACCENT),
    ("19:19", "🏠 家庭中枢", "门锁已解锁，玄关灯已开启，客厅温度26°C。欢迎回家。", SUCCESS),
]

y = Inches(1.8)
for time, device, desc, color in events:
    add_shape_bg(slide, Inches(1.0), y, Inches(0.8), Inches(0.8), color)
    tf = slide.shapes.add_textbox(Inches(1.0), y, Inches(0.8), Inches(0.8)).text_frame
    tf.paragraphs[0].text = time
    tf.paragraphs[0].font.size = Pt(12)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    add_text_box(slide, Inches(2.0), y + Inches(0.05), Inches(2.5), Inches(0.4),
                 device, font_size=16, color=color, bold=True)
    add_text_box(slide, Inches(2.0), y + Inches(0.5), Inches(10), Inches(0.5),
                 desc, font_size=14, color=BLACK)
    y += Inches(1.3)
add_slide_number(slide, 13)

# Slide 14: Four Paradigm Shifts
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, LIGHT_BG)
add_title_bar(slide, "四个范式转变", "从传统模式到智能体模式")

shifts = [
    ("从人控制到意图驱动", "你说下班了，不是指令清单\n人类从操作员变成审批者", PRIMARY),
    ("从孤岛设备到协作网络", "设备之间用群聊协议通信\n不需要人当传话筒", ACCENT),
    ("从固定规则到动态优化", "根据路况、电价、习惯\n实时计算最优解", SUCCESS),
    ("从被动响应到主动汇报", "智能体先汇报方案\n你确认后执行", ORANGE),
]

x = Inches(0.5)
for title, desc, color in shifts:
    card = add_shape_bg(slide, x, Inches(1.8), Inches(3.0), Inches(4.8), WHITE)
    add_shape_bg(slide, x, Inches(1.8), Inches(3.0), Inches(0.08), color)
    add_text_box(slide, x + Inches(0.15), Inches(2.1), Inches(2.7), Inches(1.2),
                 title, font_size=18, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.2), Inches(3.5), Inches(2.6), Inches(2.5),
                 desc, font_size=16, color=GRAY, alignment=PP_ALIGN.CENTER, line_spacing=1.6)
    x += Inches(3.2)
add_slide_number(slide, 14)

# Slide 15: Scenario 2 - Lifecycle
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, LIGHT_BG)
add_title_bar(slide, "场景二：设备全生命周期", "七个阶段，数据连续")

stages = ["购买", "配送", "安装", "日常使用", "保养维修", "二手交易", "回收报废"]
colors = [PRIMARY, ACCENT, SUCCESS, PRIMARY, ACCENT, SUCCESS, ORANGE]

x = Inches(0.3)
for i, (stage, color) in enumerate(zip(stages, colors)):
    shape = add_shape_bg(slide, x, Inches(2.0), Inches(1.5), Inches(1.0), color)
    tf = shape.text_frame
    tf.paragraphs[0].text = stage
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    if i < len(stages) - 1:
        add_text_box(slide, x + Inches(1.5), Inches(2.2), Inches(0.5), Inches(0.5),
                     "\u2192", font_size=20, color=GRAY, alignment=PP_ALIGN.CENTER)
        x += Inches(1.8)
    else:
        x += Inches(1.5)

add_text_box(slide, Inches(1.0), Inches(3.5), Inches(11.3), Inches(1.0),
             "群聊智能体贯穿全程，每个阶段设备身份不变，数据连续",
             font_size=20, color=PRIMARY, bold=True, alignment=PP_ALIGN.CENTER)
add_slide_number(slide, 15)

# Slide 16: Purchase Decision
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, LIGHT_BG)
add_title_bar(slide, "第一阶段：购买决策", "旧冰箱智能体主动汇报")

add_shape_bg(slide, Inches(1.0), Inches(1.8), Inches(5.5), Inches(4.5), WHITE)
add_text_box(slide, Inches(1.3), Inches(1.9), Inches(4.9), Inches(0.4),
             "🧊 旧冰箱智能体汇报：", font_size=18, color=PRIMARY, bold=True)
add_text_box(slide, Inches(1.3), Inches(2.5), Inches(4.9), Inches(3.5),
             "压缩机运行：43,200小时\n制冷效率：下降至65%\n维修成本约800元\n修后预计寿命<2年\n\n建议：更换更划算",
             font_size=15, color=BLACK, line_spacing=1.6)

add_shape_bg(slide, Inches(6.8), Inches(1.8), Inches(5.5), Inches(4.5), WHITE)
add_text_box(slide, Inches(7.1), Inches(1.9), Inches(4.9), Inches(0.4),
             "🛒 购物智能体推荐：", font_size=18, color=ACCENT, bold=True)
add_text_box(slide, Inches(7.1), Inches(2.5), Inches(4.9), Inches(3.5),
             "海尔 BCD-406：¥3,899\n美的 BCD-424：¥4,299\n西门子 KA92：¥5,199\n\n三平台比价完成\n京东当前最低价",
             font_size=15, color=BLACK, line_spacing=1.6)
add_slide_number(slide, 16)

# Slide 17: Installation
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, LIGHT_BG)
add_title_bar(slide, "第三阶段：安装与身份注册", "新设备上线，旧设备下线")

# New fridge
add_shape_bg(slide, Inches(1.0), Inches(1.8), Inches(5.5), Inches(4.5), WHITE)
add_text_box(slide, Inches(1.3), Inches(1.9), Inches(4.9), Inches(0.4),
             "🆕 新冰箱智能体首次上线：", font_size=18, color=SUCCESS, bold=True)
add_text_box(slide, Inches(1.3), Inches(2.5), Inches(4.9), Inches(3.5),
             "序列号：HR-2026-88A3F2\n固件版本：v3.2.1\n保修期：至2029年5月23日\n已接入家庭WiFi（厨房）",
             font_size=15, color=BLACK, line_spacing=1.6)

# Old fridge
add_shape_bg(slide, Inches(6.8), Inches(1.8), Inches(5.5), Inches(4.5), WHITE)
add_text_box(slide, Inches(7.1), Inches(1.9), Inches(4.9), Inches(0.4),
             "🔄 旧冰箱智能体最后发言：", font_size=18, color=ORANGE, bold=True)
add_text_box(slide, Inches(7.1), Inches(2.5), Inches(4.9), Inches(3.5),
             "总运行时长：43,200小时\n总耗电量：8,640度\n数据已上传至区块链存证\n\n再见。👋",
             font_size=15, color=BLACK, line_spacing=1.6)
add_slide_number(slide, 17)

# Slide 18: Daily Use
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, LIGHT_BG)
add_title_bar(slide, "第四阶段：日常使用——持续记录", "每日自动归档，每月健康报告")

# JSON
add_shape_bg(slide, Inches(1.0), Inches(1.8), Inches(6.0), Inches(4.5), RGBColor(0x2D, 0x2D, 0x2D))
add_text_box(slide, Inches(1.3), Inches(2.0), Inches(5.4), Inches(4.0),
             '{\n  "deviceId": "Fridge-002",\n  "date": "2026-05-23",\n  "metrics": {\n    "运行时长": "18.5小时",\n    "耗电量": "0.8度",\n    "平均温度": "冷藏3.2°C",\n    "开门次数": "14次",\n    "异常事件": "0次"\n  },\n  "health": "良好",\n  "efficiency": "98%"\n}',
             font_size=13, color=ACCENT, font_name='Courier New', line_spacing=1.4)

# Monthly report
add_shape_bg(slide, Inches(7.3), Inches(1.8), Inches(5.0), Inches(4.5), WHITE)
add_text_box(slide, Inches(7.6), Inches(1.9), Inches(4.4), Inches(0.4),
             "🏠 家庭中枢每月报告：", font_size=18, color=PRIMARY, bold=True)
add_text_box(slide, Inches(7.6), Inches(2.5), Inches(4.4), Inches(3.5),
             "冰箱：运行正常，效率98%\n空调：滤网需清洗\n洗衣机：脱水震动略大\n热水器：镁棒寿命剩余30%\n\n本月总电费：¥286\n（较上月下降12%）",
             font_size=15, color=BLACK, line_spacing=1.6)
add_slide_number(slide, 18)

# Slide 19: Maintenance
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, LIGHT_BG)
add_title_bar(slide, "第五阶段：保养与维修", "智能体自检 + 售后联动")

add_shape_bg(slide, Inches(1.0), Inches(1.8), Inches(11.3), Inches(4.5), WHITE)

lines = [
    "主人：冰箱最近有点异响。",
    "",
    "🧊 冰箱智能体自检：",
    "  - 压缩机：正常",
    "  - 风扇：轻微偏心振动（频率23Hz）",
    "  - 诊断：风扇轴承磨损，保修期内免费更换",
    "",
    "📞 售后智能体：",
    "  - 已预约：3月15日 10:00-10:30",
    "  - 工程师：李工（工号HR-TEC-4421）",
    "  - 工程师已收到设备诊断报告，会携带对应配件",
    "",
    "✅ 维修完成：修复后噪音从42dB降至28dB",
]
add_text_box(slide, Inches(1.5), Inches(2.0), Inches(10.3), Inches(4.0),
             "\n".join(lines), font_size=16, color=BLACK, line_spacing=1.5)
add_slide_number(slide, 19)

# Slide 20: Second-hand Trading
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, LIGHT_BG)
add_title_bar(slide, "第六阶段：二手交易", "核心价值：全生命周期数据可追溯")

add_shape_bg(slide, Inches(1.0), Inches(1.8), Inches(11.3), Inches(4.5), WHITE)

add_text_box(slide, Inches(1.5), Inches(2.0), Inches(10.3), Inches(0.4),
             "🧊 冰箱智能体生成二手车式完整报告：", font_size=20, color=PRIMARY, bold=True)

lines = [
    "型号：海尔 BCD-406 十字对开",
    "使用时长：2年1个月 | 总运行时长：21,600小时",
    "压缩机寿命：剩余75%（优秀）",
    "制冷效率：95%（正常衰减）",
    "维修记录：1次（2027年3月风扇轴承更换，保修内免费）",
    "日均开门：12次（轻度使用）",
    "",
    "当前估值：¥2,100-2,500",
    "区块链存证：完整运行数据可查，买家可验证",
]
add_text_box(slide, Inches(1.5), Inches(2.6), Inches(10.3), Inches(3.5),
             "\n".join(lines), font_size=16, color=BLACK, line_spacing=1.5)
add_slide_number(slide, 20)

# Slide 21: Data Chain
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, LIGHT_BG)
add_title_bar(slide, "全生命周期数据链", "每台设备一个数字身份证，数据不可篡改")

add_shape_bg(slide, Inches(1.0), Inches(1.8), Inches(11.3), Inches(4.5), WHITE)

lines = [
    "设备ID: HR-2026-88A3F2",
    "",
    "📦 购买记录（时间/价格/渠道/保修）",
    "📈 运行数据（每日自动记录）",
    "🔧 维修记录（时间/内容/费用）",
    "📊 健康评估（每月自动生成）",
    "💰 交易记录（二手出售/数据所有权转移）",
    "",
    "数据不可篡改，支撑二手交易信任",
]
add_text_box(slide, Inches(1.5), Inches(2.0), Inches(10.3), Inches(4.0),
             "\n".join(lines), font_size=18, color=BLACK, line_spacing=1.6)
add_slide_number(slide, 21)

# Slide 22: Technical Architecture
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, LIGHT_BG)
add_title_bar(slide, "技术架构", "系统架构总览")

layers = [
    ("自然语言层", "下班了，回家。", PRIMARY, Inches(1.8)),
    ("意图理解智能体（LLM）", "解析：回家意图 + 触发归家流程", ACCENT, Inches(3.0)),
    ("家庭中枢智能体（Orchestrator）", "协调各设备智能体 + 生成执行方案", SUCCESS, Inches(4.2)),
]

for title, desc, color, y in layers:
    shape = add_shape_bg(slide, Inches(3.0), y, Inches(7.3), Inches(0.9), color)
    tf = shape.text_frame
    tf.paragraphs[0].text = title
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(13)
    p2.font.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
    p2.alignment = PP_ALIGN.CENTER
    
    if y < Inches(4.2):
        add_text_box(slide, Inches(6.3), y + Inches(0.9), Inches(0.7), Inches(0.5),
                     "\u25BC", font_size=20, color=GRAY, alignment=PP_ALIGN.CENTER)

devices = ["汽车", "空调", "冰箱", "门窗", "灯光"]
x = Inches(1.0)
for dev in devices:
    shape = add_shape_bg(slide, x, Inches(5.3), Inches(2.0), Inches(0.7), WHITE)
    tf = shape.text_frame
    tf.paragraphs[0].text = dev
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.color.rgb = BLACK
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    x += Inches(2.2)

add_shape_bg(slide, Inches(1.5), Inches(6.3), Inches(10.3), Inches(0.7), DARK_BG)
add_text_box(slide, Inches(1.5), Inches(6.4), Inches(10.3), Inches(0.5),
             "群聊协议（钉钉/微信/自建）— 设备间通信 + 状态广播 + 方案审批",
             font_size=14, color=WHITE, alignment=PP_ALIGN.CENTER)
add_slide_number(slide, 22)

# Slide 23: Key Technologies
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, LIGHT_BG)
add_title_bar(slide, "关键技术支撑", "七大技术栈")

tech_items = [
    ("LLM意图理解", "自然语言指令解析", PRIMARY),
    ("多智能体协作", "设备之间自主协商", ACCENT),
    ("群聊协议", "设备间通信统一接口", SUCCESS),
    ("区块链存证", "数据不可篡改", ORANGE),
    ("预测性维护", "基于数据预测故障", PRIMARY),
    ("数字身份管理", "设备唯一ID", ACCENT),
    ("边缘计算", "本地快速响应", SUCCESS),
]

x = Inches(0.5)
y = Inches(1.8)
for title, desc, color in tech_items:
    card = add_shape_bg(slide, x, y, Inches(3.5), Inches(1.5), WHITE)
    add_shape_bg(slide, x, y, Inches(3.5), Inches(0.06), color)
    add_text_box(slide, x + Inches(0.15), y + Inches(0.15), Inches(3.2), Inches(0.4),
                 title, font_size=18, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.15), y + Inches(0.6), Inches(3.2), Inches(0.6),
                 desc, font_size=14, color=GRAY, alignment=PP_ALIGN.CENTER)
    x += Inches(3.8)
    if x > Inches(10):
        x = Inches(0.5)
        y += Inches(1.8)
add_slide_number(slide, 23)

# Slide 24: Key Challenges
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, LIGHT_BG)
add_title_bar(slide, "关键挑战", "不能只唱赞歌")

challenges = [
    ("身份安全", "如何确认群里的空调智能体是真的空调，不是钓鱼？", PRIMARY),
    ("责任归属", "智能体误操作导致损失，谁负责？人、厂商、平台？", ACCENT),
    ("隐私边界", "传感器持续感知环境 = 持续收集数据，如何保障？", SUCCESS),
    ("数字鸿沟", "不会用智能体的老年人/弱势群体被进一步边缘化？", ORANGE),
    ("生态碎片化", "小米智能体能跟华为智能体对话吗？标准谁定？", PRIMARY),
]

y = Inches(1.8)
for idx, (title, desc, color) in enumerate(challenges):
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.0), y + Inches(0.1), Inches(0.5), Inches(0.5))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    tf = circle.text_frame
    tf.paragraphs[0].text = str(idx + 1)
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    add_text_box(slide, Inches(1.8), y + Inches(0.05), Inches(3), Inches(0.4),
                 title, font_size=20, color=color, bold=True)
    add_text_box(slide, Inches(1.8), y + Inches(0.5), Inches(10), Inches(0.5),
                 desc, font_size=16, color=GRAY)
    y += Inches(1.1)
add_slide_number(slide, 24)

# Slide 25: National Strategy
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, LIGHT_BG)
add_title_bar(slide, "国家战略与教育对齐", "国务院2027年智能体普及率70%目标")

nets = ["算力网", "电力网", "物流网", "通信网", "数据网", "安全网"]
x = Inches(0.8)
for net in nets:
    shape = add_shape_bg(slide, x, Inches(2.0), Inches(1.8), Inches(0.8), PRIMARY)
    tf = shape.text_frame
    tf.paragraphs[0].text = net
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    x += Inches(2.0)

add_shape_bg(slide, Inches(1.0), Inches(3.5), Inches(11.3), Inches(2.5), PRIMARY)
add_text_box(slide, Inches(1.5), Inches(3.8), Inches(10.3), Inches(0.6),
             "高校不能等产业成熟了再教，", font_size=24, color=WHITE, bold=True,
             alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.5), Inches(4.4), Inches(10.3), Inches(0.6),
             "必须在基础设施建设期就培养对应人才。", font_size=24, color=ACCENT, bold=True,
             alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.5), Inches(5.2), Inches(10.3), Inches(0.5),
             "培养智能体原生代人才，而非转型代", font_size=18, color=RGBColor(0xE0, 0xE0, 0xE0),
             alignment=PP_ALIGN.CENTER)
add_slide_number(slide, 25)

# Slide 26: ZJGU Opportunity
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, LIGHT_BG)
add_title_bar(slide, "浙江工商大学的机会", "战略位置")

add_shape_bg(slide, Inches(1.0), Inches(1.8), Inches(11.3), Inches(4.5), WHITE)

add_text_box(slide, Inches(1.5), Inches(2.0), Inches(10.3), Inches(0.4),
             "战略位置：", font_size=22, color=PRIMARY, bold=True)

lines = [
    "杭州：数字经济第一城",
    "钉钉总部：智能体生态天然试验场",
    "产业密集：跨境电商、智能制造、智慧物流",
    "",
    "培养智能体原生代人才，而非转型代",
]
add_text_box(slide, Inches(1.5), Inches(2.6), Inches(10.3), Inches(3.5),
             "\n".join(lines), font_size=18, color=BLACK, line_spacing=1.6)
add_slide_number(slide, 26)

# Slide 27: Discussion
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, LIGHT_BG)
add_title_bar(slide, "课堂讨论", "6个讨论题")

questions = [
    "数字员工和传统软件工具有什么本质区别？",
    "如果导航智能体判断错误导致迟到，责任归谁？",
    "设计一个离家模式，安排哪些智能体协作？",
    "如何平衡便利性和隐私保护？",
    "70%普及率目标，最大的障碍是什么？",
    "智能家居智能体需要哪些学科的知识？",
]

y = Inches(1.8)
for idx, q in enumerate(questions):
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.0), y + Inches(0.1), Inches(0.5), Inches(0.5))
    circle.fill.solid()
    circle.fill.fore_color.rgb = PRIMARY
    circle.line.fill.background()
    tf = circle.text_frame
    tf.paragraphs[0].text = str(idx + 1)
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    add_text_box(slide, Inches(1.8), y + Inches(0.05), Inches(10), Inches(0.5),
                 q, font_size=18, color=BLACK)
    y += Inches(0.9)
add_slide_number(slide, 27)

# Slide 28: Projects
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, LIGHT_BG)
add_title_bar(slide, "项目实践建议", "5个实践项目")

headers = ["项目", "难度", "技术栈"]
rows = [
    ["智能家居群聊模拟器", "⭐⭐", "Python + 群聊API"],
    ["设备数字身份系统", "⭐⭐⭐", "区块链 + Python"],
    ["预测性维护算法", "⭐⭐⭐", "机器学习"],
    ["二手设备估值模型", "⭐⭐⭐", "数据分析"],
    ["多智能体调度优化", "⭐⭐⭐⭐", "运筹优化"],
]
add_table(slide, Inches(1.0), Inches(1.8), Inches(11.3), Inches(3.5), headers, rows)

add_text_box(slide, Inches(1.0), Inches(5.5), Inches(11.3), Inches(0.4),
             "课时安排：4-8课时/项目，可根据学生水平调整",
             font_size=16, color=GRAY, alignment=PP_ALIGN.CENTER)
add_slide_number(slide, 28)

# Slide 29: Summary
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_shape_bg(slide, Inches(0), Inches(3.0), Inches(13.333), Inches(0.04), ACCENT)
add_text_box(slide, Inches(1.5), Inches(1.0), Inches(10.3), Inches(0.8),
             "核心结论", font_size=36, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

keywords = [
    ("数字员工", "设备有身份，是同事不是工具", PRIMARY),
    ("群聊即OS", "对话即交互，意图即指令", ACCENT),
    ("全生命周期", "从购买到报废，数据连续可追溯", SUCCESS),
    ("国家战略", "70%普及率目标，教育必须对齐产业", ORANGE),
]

x = Inches(0.8)
for title, desc, color in keywords:
    card = add_shape_bg(slide, x, Inches(1.8), Inches(2.8), Inches(2.5), WHITE)
    add_shape_bg(slide, x, Inches(1.8), Inches(2.8), Inches(0.06), color)
    add_text_box(slide, x + Inches(0.15), Inches(2.1), Inches(2.5), Inches(0.5),
                 title, font_size=22, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.15), Inches(2.7), Inches(2.5), Inches(1.2),
                 desc, font_size=14, color=GRAY, alignment=PP_ALIGN.CENTER, line_spacing=1.5)
    x += Inches(3.1)

add_text_box(slide, Inches(1.5), Inches(5.0), Inches(10.3), Inches(1.0),
             "智能体不是又一个技术风口，\n而是人机交互的范式转移。",
             font_size=24, color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.5), Inches(6.2), Inches(10.3), Inches(0.5),
             "从人适应机器到机器适应人",
             font_size=18, color=GRAY, alignment=PP_ALIGN.CENTER)
add_slide_number(slide, 29)

# Slide 30: Q&A
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_shape_bg(slide, Inches(0), Inches(3.0), Inches(13.333), Inches(0.04), ACCENT)
add_text_box(slide, Inches(1.5), Inches(2.0), Inches(10.3), Inches(1.0),
             "谢谢！", font_size=48, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.5), Inches(3.5), Inches(10.3), Inches(0.8),
             "欢迎提问和讨论", font_size=28, color=ACCENT, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.5), Inches(5.0), Inches(10.3), Inches(0.5),
             "诸葛斌 / 虾尔 AI / 2026年5月", font_size=16, color=GRAY, alignment=PP_ALIGN.CENTER)

# Save
output_path = "/home/admin/.openclaw/workspace/智能体数字员工与智能家居_完整版.pptx"
prs.save(output_path)
print(f"PPT saved: {output_path}")
print(f"File size: {os.path.getsize(output_path) / 1024:.1f} KB")
