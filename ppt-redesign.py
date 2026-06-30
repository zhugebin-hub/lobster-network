#!/usr/bin/env python3
"""毕业论文答辩PPT重新排版 - 浅色主题"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import copy

# ========== 配色方案 ==========
BG_COLOR = RGBColor(0xFA, 0xF8, 0xF0)         # 暖白背景
ACCENT_GOLD = RGBColor(0xC4, 0x9A, 0x4B)      # 金色点缀
ACCENT_DARK = RGBColor(0x3A, 0x3A, 0x2E)      # 深褐文字
ACCENT_MID = RGBColor(0x6B, 0x5B, 0x3E)       # 中褐
ACCENT_LIGHT = RGBColor(0xE8, 0xE0, 0xCC)     # 浅金边框
CARD_BG = RGBColor(0xFF, 0xFF, 0xF5)          # 卡片背景
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LINE_COLOR = RGBColor(0xD4, 0xCA, 0xA0)       # 分割线

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height

# ========== 工具函数 ==========

def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, fill_color, border_color=None, border_width=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(border_width or 1)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    # 圆角调整
    shape.adjustments[0] = 0.03
    return shape

def add_rect(slide, left, top, width, height, fill_color, border_color=None, border_width=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(border_width or 1)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def add_textbox(slide, left, top, width, height, text, font_size=18, bold=False, color=ACCENT_DARK, alignment=PP_ALIGN.LEFT, font_name='微软雅黑'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_para(tf, text, font_size=18, bold=False, color=ACCENT_DARK, alignment=PP_ALIGN.LEFT, space_before=Pt(4), space_after=Pt(2), font_name='微软雅黑'):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    if space_before:
        p.space_before = space_before
    if space_after:
        p.space_after = space_after
    return p

def add_decorative_line(slide, left, top, width, color=ACCENT_GOLD, height=Pt(2)):
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()
    line.shadow.inherit = False
    return line

def add_top_bar(slide):
    """顶部装饰条"""
    add_decorative_line(slide, Inches(0), Inches(0), SLIDE_W, ACCENT_GOLD, Pt(3))

def add_bottom_bar(slide):
    """底部装饰条"""
    add_decorative_line(slide, Inches(0.8), Inches(7.1), Inches(11.7), LINE_COLOR, Pt(1))

def add_slide_number(slide, num, total=11):
    add_textbox(slide, Inches(11.5), Inches(7.15), Inches(1.5), Inches(0.3),
                f"{num} / {total}", font_size=10, color=RGBColor(0xAA, 0xAA, 0xAA),
                alignment=PP_ALIGN.RIGHT)

def add_chapter_label(slide, num, title):
    """章节标签 - 左上角"""
    label = add_shape(slide, Inches(0.6), Inches(0.5), Inches(3.5), Inches(0.55),
                      ACCENT_GOLD, border_color=None)
    label.text_frame.paragraphs[0].text = f"第{num}部分  {title}"
    label.text_frame.paragraphs[0].font.size = Pt(14)
    label.text_frame.paragraphs[0].font.bold = True
    label.text_frame.paragraphs[0].font.color.rgb = WHITE
    label.text_frame.paragraphs[0].font.name = '微软雅黑'
    label.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
    label.text_frame.margin_left = Inches(0.15)
    label.text_frame.margin_top = Inches(0.08)

# ============================================================
# 第1页：封面
# ============================================================
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide1, BG_COLOR)
add_top_bar(slide1)

# 装饰：左侧金色竖条
add_rect(slide1, Inches(0.6), Inches(1.5), Pt(4), Inches(4.5), ACCENT_GOLD)

# 主标题区域
add_textbox(slide1, Inches(1.2), Inches(1.6), Inches(10), Inches(0.7),
            "毕业论文答辩", font_size=22, bold=True, color=ACCENT_GOLD, font_name='微软雅黑')

add_textbox(slide1, Inches(1.2), Inches(2.5), Inches(10), Inches(1.2),
            "佛教中国化视域下\n十善业道与儒家五常的融合路径研究",
            font_size=36, bold=True, color=ACCENT_DARK, font_name='微软雅黑')

# 分割线
add_decorative_line(slide1, Inches(1.2), Inches(4.2), Inches(3), ACCENT_GOLD, Pt(2.5))

# 信息
tb = add_textbox(slide1, Inches(1.2), Inches(4.6), Inches(6), Inches(2),
                 "", font_size=16, color=ACCENT_MID)
tf = tb.text_frame
tf.paragraphs[0].text = "答辩人：释果顺"
tf.paragraphs[0].font.size = Pt(18)
tf.paragraphs[0].font.color.rgb = ACCENT_MID
tf.paragraphs[0].font.name = '微软雅黑'
add_para(tf, "指导教师：郑根成 教授", font_size=18, color=ACCENT_MID)
add_para(tf, "研究方向：佛教伦理", font_size=18, color=ACCENT_MID)
add_para(tf, "浙江省宗教界「双通」人才研修班", font_size=16, color=RGBColor(0x99, 0x88, 0x66))
add_para(tf, "2026年6月", font_size=16, color=RGBColor(0x99, 0x88, 0x66))

# 右下角装饰
add_rect(slide1, Inches(11.5), Inches(6.5), Inches(1.2), Inches(0.06), ACCENT_GOLD)

add_slide_number(slide1, 1)

# ============================================================
# 第2页：汇报框架
# ============================================================
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide2, BG_COLOR)
add_top_bar(slide2)
add_chapter_label(slide2, "二", "汇报框架")

add_textbox(slide2, Inches(0.6), Inches(1.3), Inches(10), Inches(0.6),
            "汇报框架", font_size=32, bold=True, color=ACCENT_DARK)

items = [
    ("01", "选题背景与研究问题"),
    ("02", "研究框架与核心内容"),
    ("03", "各章要点概述"),
    ("04", "主要创新"),
    ("05", "不足与展望"),
]

card_w = Inches(2.1)
card_h = Inches(1.6)
start_x = Inches(0.8)
gap = Inches(0.35)

for i, (num, title) in enumerate(items):
    x = start_x + i * (card_w + gap)
    y = Inches(2.4)
    
    # 卡片
    card = add_shape(slide2, x, y, card_w, card_h, CARD_BG, LINE_COLOR, 1)
    
    # 编号
    num_shape = add_shape(slide2, x + Inches(0.1), y + Inches(0.15), Inches(0.5), Inches(0.5),
                          ACCENT_GOLD)
    num_shape.text_frame.paragraphs[0].text = num
    num_shape.text_frame.paragraphs[0].font.size = Pt(20)
    num_shape.text_frame.paragraphs[0].font.bold = True
    num_shape.text_frame.paragraphs[0].font.color.rgb = WHITE
    num_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    num_shape.text_frame.margin_top = Inches(0.05)
    
    # 标题
    add_textbox(slide2, x + Inches(0.15), y + Inches(0.85), card_w - Inches(0.3), Inches(0.6),
                title, font_size=15, bold=True, color=ACCENT_DARK, alignment=PP_ALIGN.CENTER)

add_bottom_bar(slide2)
add_slide_number(slide2, 2)

# ============================================================
# 第3页：选题背景与研究问题
# ============================================================
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide3, BG_COLOR)
add_top_bar(slide3)
add_chapter_label(slide3, "三", "选题背景")

add_textbox(slide3, Inches(0.6), Inches(1.3), Inches(10), Inches(0.6),
            "选题背景与研究问题", font_size=32, bold=True, color=ACCENT_DARK)

# 左侧：研究背景
card1 = add_shape(slide3, Inches(0.6), Inches(2.2), Inches(5.8), Inches(4.5), CARD_BG, LINE_COLOR, 1.5)
tb1 = add_textbox(slide3, Inches(0.9), Inches(2.35), Inches(5.2), Inches(0.5),
                  "▎ 研究背景", font_size=18, bold=True, color=ACCENT_GOLD)

bg_items = [
    "佛教中国化研究中，伦理融合是最核心的维度之一",
    "现有成果多聚焦「五戒」与「五常」的静态格义比附",
    "五戒偏重外在底线禁戒，难以契合心性修养诉求",
]
tb1b = add_textbox(slide3, Inches(1.0), Inches(3.0), Inches(5.0), Inches(3.5), "", font_size=15, color=ACCENT_MID)
tf1 = tb1b.text_frame
tf1.word_wrap = True
for idx, item in enumerate(bg_items):
    if idx == 0:
        tf1.paragraphs[0].text = f"  •  {item}"
        tf1.paragraphs[0].font.size = Pt(15)
        tf1.paragraphs[0].font.color.rgb = ACCENT_MID
        tf1.paragraphs[0].font.name = '微软雅黑'
    else:
        p = add_para(tf1, f"  •  {item}", font_size=15, color=ACCENT_MID, space_before=Pt(10))

# 右侧：核心问题
card2 = add_shape(slide3, Inches(6.8), Inches(2.2), Inches(5.8), Inches(4.5), CARD_BG, LINE_COLOR, 1.5)
tb2 = add_textbox(slide3, Inches(7.1), Inches(2.35), Inches(5.2), Inches(0.5),
                  "▎ 核心问题", font_size=18, bold=True, color=ACCENT_GOLD)

q_items = [
    "佛教如何通过「十善」与儒家「五常」的概念互释与心性会通，实现与中华主流价值的深度嵌合？",
    "这一历史经验对当代佛教中国化有何启示？",
]
tb2b = add_textbox(slide3, Inches(7.1), Inches(3.0), Inches(5.3), Inches(3.5), "", font_size=15, color=ACCENT_MID)
tf2 = tb2b.text_frame
tf2.word_wrap = True
for idx, item in enumerate(q_items):
    if idx == 0:
        tf2.paragraphs[0].text = f"  ❖  {item}"
        tf2.paragraphs[0].font.size = Pt(15)
        tf2.paragraphs[0].font.color.rgb = ACCENT_MID
        tf2.paragraphs[0].font.name = '微软雅黑'
    else:
        p = add_para(tf2, f"  ❖  {item}", font_size=15, color=ACCENT_MID, space_before=Pt(14))

add_bottom_bar(slide3)
add_slide_number(slide3, 3)

# ============================================================
# 第4页：研究框架（三维分析模型）
# ============================================================
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide4, BG_COLOR)
add_top_bar(slide4)
add_chapter_label(slide4, "四", "研究框架")

add_textbox(slide4, Inches(0.6), Inches(1.3), Inches(10), Inches(0.6),
            "研究框架：三维分析模型", font_size=32, bold=True, color=ACCENT_DARK)

# 三维模型标签
model_label = add_shape(slide4, Inches(4.2), Inches(2.1), Inches(4.8), Inches(0.6), ACCENT_GOLD)
model_label.text_frame.paragraphs[0].text = "「学理建构 → 制度运作 → 现代转化」"
model_label.text_frame.paragraphs[0].font.size = Pt(18)
model_label.text_frame.paragraphs[0].font.bold = True
model_label.text_frame.paragraphs[0].font.color.rgb = WHITE
model_label.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
model_label.text_frame.margin_top = Inches(0.1)

# 四章卡片
chapters = [
    ("第一章", "理论渊源", "十善与五常的\n「异质同构」", "01"),
    ("第二章", "历史演进", "魏晋格义 → 隋唐\n判教 → 宋明圆融", "02"),
    ("第三章", "深层机制", "义理同构 /\n心性贯通 /\n行为规训", "03"),
    ("第四章", "现代转化", "经济伦理 /\n生态伦理 /\n网络伦理", "04"),
]

cw = Inches(2.7)
ch = Inches(3.2)
sx = Inches(0.6)
sg = Inches(0.3)

for i, (ch_name, ch_sub, ch_desc, ch_num) in enumerate(chapters):
    x = sx + i * (cw + sg)
    y = Inches(3.0)
    
    card = add_shape(slide4, x, y, cw, ch, CARD_BG, LINE_COLOR, 1.5)
    
    # 编号
    nshape = add_shape(slide4, x + Inches(0.1), y + Inches(0.12), Inches(0.45), Inches(0.45), ACCENT_GOLD)
    nshape.text_frame.paragraphs[0].text = ch_num
    nshape.text_frame.paragraphs[0].font.size = Pt(16)
    nshape.text_frame.paragraphs[0].font.bold = True
    nshape.text_frame.paragraphs[0].font.color.rgb = WHITE
    nshape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    nshape.text_frame.margin_top = Inches(0.06)
    
    # 章名
    add_textbox(slide4, x + Inches(0.6), y + Inches(0.15), cw - Inches(0.8), Inches(0.4),
                ch_name, font_size=15, bold=True, color=ACCENT_GOLD)
    
    # 副标题
    add_textbox(slide4, x + Inches(0.15), y + Inches(0.65), cw - Inches(0.3), Inches(0.4),
                ch_sub, font_size=17, bold=True, color=ACCENT_DARK, alignment=PP_ALIGN.CENTER)
    
    # 描述
    add_textbox(slide4, x + Inches(0.2), y + Inches(1.3), cw - Inches(0.4), Inches(1.8),
                ch_desc, font_size=14, color=ACCENT_MID, alignment=PP_ALIGN.CENTER)

add_bottom_bar(slide4)
add_slide_number(slide4, 4)

# ============================================================
# 第5页：第一章 - 理论渊源
# ============================================================
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide5, BG_COLOR)
add_top_bar(slide5)
add_chapter_label(slide5, "五", "理论渊源")

add_textbox(slide5, Inches(0.6), Inches(1.3), Inches(10), Inches(0.6),
            "第一章  理论渊源：十善业道与五常", font_size=30, bold=True, color=ACCENT_DARK)

# 核心概念卡片
card_left = add_shape(slide5, Inches(0.6), Inches(2.2), Inches(5.8), Inches(4.5), CARD_BG, LINE_COLOR, 1.5)
add_textbox(slide5, Inches(0.9), Inches(2.35), Inches(5), Inches(0.4),
            "▎ 核心概念：异质同构", font_size=18, bold=True, color=ACCENT_GOLD)

tb5a = add_textbox(slide5, Inches(1.0), Inches(2.9), Inches(5.2), Inches(3.5), "", font_size=15, color=ACCENT_MID)
tf5a = tb5a.text_frame
tf5a.word_wrap = True
tf5a.paragraphs[0].text = "  •  终极目标不同：出世解脱 vs 入世治平"
tf5a.paragraphs[0].font.size = Pt(15)
tf5a.paragraphs[0].font.color.rgb = ACCENT_MID
tf5a.paragraphs[0].font.name = '微软雅黑'
add_para(tf5a, "  •  心性修养路径：结构性共鸣", font_size=15, color=ACCENT_MID, space_before=Pt(12))

# 右侧：十善 > 五戒
card_right = add_shape(slide5, Inches(6.8), Inches(2.2), Inches(5.8), Inches(4.5), CARD_BG, LINE_COLOR, 1.5)
add_textbox(slide5, Inches(7.1), Inches(2.35), Inches(5), Inches(0.4),
            "▎ 十善 > 五戒：关键突破", font_size=18, bold=True, color=ACCENT_GOLD)

# 三业分类
ye_items = [
    ("身业", "不杀、不盗、不邪淫"),
    ("口业", "不妄语、不两舌、不恶口、不绮语"),
    ("意业", "不贪、不嗔、不邪见  ← 关键突破"),
]

ty = Inches(3.0)
for label, desc in ye_items:
    # 标签
    lbl = add_shape(slide5, Inches(7.1), ty, Inches(0.7), Inches(0.4), ACCENT_GOLD)
    lbl.text_frame.paragraphs[0].text = label
    lbl.text_frame.paragraphs[0].font.size = Pt(13)
    lbl.text_frame.paragraphs[0].font.bold = True
    lbl.text_frame.paragraphs[0].font.color.rgb = WHITE
    lbl.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    lbl.text_frame.margin_top = Inches(0.04)
    
    # 描述
    add_textbox(slide5, Inches(7.95), ty + Inches(0.02), Inches(4.5), Inches(0.4),
                desc, font_size=15, color=ACCENT_MID)
    ty += Inches(0.55)

# 底部结论
add_textbox(slide5, Inches(7.1), Inches(5.0), Inches(5.2), Inches(0.6),
            "→ 与儒家「诚意正心」实现本体论对接",
            font_size=16, bold=True, color=ACCENT_GOLD)

add_bottom_bar(slide5)
add_slide_number(slide5, 5)

# ============================================================
# 第6页：第二章 - 历史演进
# ============================================================
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide6, BG_COLOR)
add_top_bar(slide6)
add_chapter_label(slide6, "六", "历史演进")

add_textbox(slide6, Inches(0.6), Inches(1.3), Inches(10), Inches(0.6),
            "第二章  历史演进：从格义到圆融", font_size=30, bold=True, color=ACCENT_DARK)

# 时间线
timeline_y = Inches(3.0)
line_w = Inches(11.5)
add_decorative_line(slide6, Inches(0.8), timeline_y + Inches(0.25), line_w, LINE_COLOR, Pt(2))

periods = [
    ("魏晋", "格义比附\n五戒配五常", "被动防御"),
    ("隋唐", "判教统摄\n宗派伦理体系化", "体系建构"),
    ("宋明", "深度圆融\n契嵩《辅教编》\n功过格、乡约", "深度嵌合"),
    ("近现代", "人间佛教\n范式转换", "现代转型"),
]

pw = Inches(2.6)
pgap = Inches(0.35)

for i, (period, desc, tag) in enumerate(periods):
    x = Inches(0.8) + i * (pw + pgap)
    
    # 圆点
    dot = slide6.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.9), timeline_y, Inches(0.35), Inches(0.35))
    dot.fill.solid()
    dot.fill.fore_color.rgb = ACCENT_GOLD
    dot.line.fill.background()
    
    # 时期标签
    period_card = add_shape(slide6, x, Inches(3.6), pw, Inches(0.55), ACCENT_GOLD)
    period_card.text_frame.paragraphs[0].text = period
    period_card.text_frame.paragraphs[0].font.size = Pt(17)
    period_card.text_frame.paragraphs[0].font.bold = True
    period_card.text_frame.paragraphs[0].font.color.rgb = WHITE
    period_card.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    period_card.text_frame.margin_top = Inches(0.08)
    
    # 描述
    add_textbox(slide6, x + Inches(0.1), Inches(4.35), pw - Inches(0.2), Inches(2.0),
                desc, font_size=13, color=ACCENT_MID, alignment=PP_ALIGN.CENTER)
    
    # 阶段标签
    tag_card = add_shape(slide6, x + Inches(0.4), Inches(6.5), Inches(1.8), Inches(0.4), CARD_BG, ACCENT_GOLD, 1)
    tag_card.text_frame.paragraphs[0].text = tag
    tag_card.text_frame.paragraphs[0].font.size = Pt(12)
    tag_card.text_frame.paragraphs[0].font.bold = True
    tag_card.text_frame.paragraphs[0].font.color.rgb = ACCENT_GOLD
    tag_card.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    tag_card.text_frame.margin_top = Inches(0.06)

add_bottom_bar(slide6)
add_slide_number(slide6, 6)

# ============================================================
# 第7页：第三章 - 融合的内在机制
# ============================================================
slide7 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide7, BG_COLOR)
add_top_bar(slide7)
add_chapter_label(slide7, "七", "内在机制")

add_textbox(slide7, Inches(0.6), Inches(1.3), Inches(10), Inches(0.6),
            "第三章  融合的内在机制", font_size=30, bold=True, color=ACCENT_DARK)

mechanisms = [
    ("义理同构机制", "概念对勘与本体论重置",
     ["「业报」→「心性」的哲学升华"],
     "01"),
    ("心性贯通机制", "",
     ["「佛性」与「良知」的心理学交融",
      "十善微观修持 ↔ 儒家格物之功"],
     "02"),
    ("行为规训机制", "",
     ["功过格：伦理指标的量化折算",
      "基层社会的秩序重塑"],
     "03"),
]

mw = Inches(3.5)
mgap = Inches(0.3)

for i, (name, subtitle, items, num) in enumerate(mechanisms):
    x = Inches(0.6) + i * (mw + mgap)
    y = Inches(2.2)
    
    card = add_shape(slide7, x, y, mw, Inches(4.5), CARD_BG, LINE_COLOR, 1.5)
    
    # 编号
    nshape = add_shape(slide7, x + Inches(0.1), y + Inches(0.12), Inches(0.4), Inches(0.4), ACCENT_GOLD)
    nshape.text_frame.paragraphs[0].text = num
    nshape.text_frame.paragraphs[0].font.size = Pt(14)
    nshape.text_frame.paragraphs[0].font.bold = True
    nshape.text_frame.paragraphs[0].font.color.rgb = WHITE
    nshape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    nshape.text_frame.margin_top = Inches(0.06)
    
    # 机制名
    add_textbox(slide7, x + Inches(0.55), y + Inches(0.15), mw - Inches(0.7), Inches(0.4),
                name, font_size=17, bold=True, color=ACCENT_GOLD)
    
    # 副标题
    if subtitle:
        add_textbox(slide7, x + Inches(0.15), y + Inches(0.65), mw - Inches(0.3), Inches(0.35),
                    subtitle, font_size=13, color=ACCENT_MID, alignment=PP_ALIGN.CENTER)
    
    # 要点
    iy = y + Inches(1.0) if not subtitle else y + Inches(1.1)
    for item in items:
        add_textbox(slide7, x + Inches(0.2), iy, mw - Inches(0.4), Inches(0.4),
                    f"  ◆  {item}", font_size=13, color=ACCENT_MID)
        iy += Inches(0.5)

# 底部箭头关系
add_textbox(slide7, Inches(2.5), Inches(6.8), Inches(8), Inches(0.4),
            "义理同构（基础） → 心性贯通（桥梁） → 行为规训（落地）",
            font_size=13, bold=True, color=ACCENT_GOLD, alignment=PP_ALIGN.CENTER)

add_bottom_bar(slide7)
add_slide_number(slide7, 7)

# ============================================================
# 第8页：第四章 - 现代转化
# ============================================================
slide8 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide8, BG_COLOR)
add_top_bar(slide8)
add_chapter_label(slide8, "八", "现代转化")

add_textbox(slide8, Inches(0.6), Inches(1.3), Inches(10), Inches(0.6),
            "第四章  现代转化与当代价值", font_size=30, bold=True, color=ACCENT_DARK)

# 三大领域
fields = [
    ("经济伦理", "契约精神重塑\n诚信建设", "📜"),
    ("生态伦理", "戒杀护生\n深层生态学", "🌿"),
    ("网络伦理", "口业四善\n数字交往规范", "🌐"),
]

fw = Inches(3.5)
fgap = Inches(0.3)

for i, (name, desc, icon) in enumerate(fields):
    x = Inches(0.6) + i * (fw + fgap)
    y = Inches(2.3)
    
    card = add_shape(slide8, x, y, fw, Inches(2.5), CARD_BG, LINE_COLOR, 1.5)
    
    # 图标
    add_textbox(slide8, x + Inches(0.15), y + Inches(0.15), Inches(0.5), Inches(0.5),
                icon, font_size=24, alignment=PP_ALIGN.CENTER)
    
    # 名称
    add_textbox(slide8, x + Inches(0.1), y + Inches(0.7), fw - Inches(0.2), Inches(0.4),
                name, font_size=18, bold=True, color=ACCENT_GOLD, alignment=PP_ALIGN.CENTER)
    
    # 描述
    add_textbox(slide8, x + Inches(0.2), y + Inches(1.2), fw - Inches(0.4), Inches(1.2),
                desc, font_size=14, color=ACCENT_MID, alignment=PP_ALIGN.CENTER)

# 实践路径
card_path = add_shape(slide8, Inches(0.6), Inches(5.1), Inches(11.2), Inches(1.5), CARD_BG, LINE_COLOR, 1.5)
add_textbox(slide8, Inches(0.9), Inches(5.2), Inches(5), Inches(0.4),
            "▎ 实践路径", font_size=16, bold=True, color=ACCENT_GOLD)

path_items = [
    "讲经示范体系  /  院校课程模块  /  社区公益服务",
    "个体心理调适：「意三善」与现代正念疗愈",
]
tb8 = add_textbox(slide8, Inches(1.0), Inches(5.7), Inches(10.5), Inches(0.8), "", font_size=14, color=ACCENT_MID)
tf8 = tb8.text_frame
tf8.word_wrap = True
tf8.paragraphs[0].text = f"  •  {path_items[0]}"
tf8.paragraphs[0].font.size = Pt(14)
tf8.paragraphs[0].font.color.rgb = ACCENT_MID
tf8.paragraphs[0].font.name = '微软雅黑'
add_para(tf8, f"  •  {path_items[1]}", font_size=14, color=ACCENT_MID, space_before=Pt(8))

add_bottom_bar(slide8)
add_slide_number(slide8, 8)

# ============================================================
# 第9页：主要创新
# ============================================================
slide9 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide9, BG_COLOR)
add_top_bar(slide9)
add_chapter_label(slide9, "九", "主要创新")

add_textbox(slide9, Inches(0.6), Inches(1.3), Inches(10), Inches(0.6),
            "主要创新", font_size=32, bold=True, color=ACCENT_DARK)

innovations = [
    ("视角创新", "突破五戒—五常静态比附范式，\n系统推进十善业道与五常的深层对勘"),
    ("概念创新", "提炼「异质同构」为统摄性\n分析范畴"),
    ("框架创新", "建构「学理—制度—现代」\n三维分析模型，贯通历史经验\n与当代应用"),
]

iw = Inches(3.5)
igap = Inches(0.3)

for i, (title, desc) in enumerate(innovations):
    x = Inches(0.6) + i * (iw + igap)
    y = Inches(2.3)
    
    card = add_shape(slide9, x, y, iw, Inches(4.0), CARD_BG, ACCENT_GOLD, 2)
    
    # 顶部色块
    top_bar = add_rect(slide9, x, y, iw, Inches(0.6), ACCENT_GOLD)
    top_bar.text_frame.paragraphs[0].text = title
    top_bar.text_frame.paragraphs[0].font.size = Pt(18)
    top_bar.text_frame.paragraphs[0].font.bold = True
    top_bar.text_frame.paragraphs[0].font.color.rgb = WHITE
    top_bar.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    top_bar.text_frame.margin_top = Inches(0.1)
    
    # 描述
    add_textbox(slide9, x + Inches(0.25), y + Inches(0.9), iw - Inches(0.5), Inches(2.8),
                desc, font_size=15, color=ACCENT_MID, alignment=PP_ALIGN.CENTER)

add_bottom_bar(slide9)
add_slide_number(slide9, 9)

# ============================================================
# 第10页：不足与展望
# ============================================================
slide10 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide10, BG_COLOR)
add_top_bar(slide10)
add_chapter_label(slide10, "十", "不足与展望")

add_textbox(slide10, Inches(0.6), Inches(1.3), Inches(10), Inches(0.6),
            "不足与展望", font_size=32, bold=True, color=ACCENT_DARK)

# 左侧：不足
card_l = add_shape(slide10, Inches(0.6), Inches(2.2), Inches(5.8), Inches(4.3), CARD_BG, LINE_COLOR, 1.5)
add_textbox(slide10, Inches(0.9), Inches(2.35), Inches(5), Inches(0.4),
            "▎ 两点不足", font_size=18, bold=True, color=RGBColor(0xCC, 0x66, 0x33))

tb10l = add_textbox(slide10, Inches(1.0), Inches(2.95), Inches(5.2), Inches(3.3), "", font_size=15, color=ACCENT_MID)
tf10l = tb10l.text_frame
tf10l.word_wrap = True
tf10l.paragraphs[0].text = "  ① 史料偏重精英规范性文本，基层碑刻、宝卷等描述性文本挖掘不足，存在「俯视感」"
tf10l.paragraphs[0].font.size = Pt(15)
tf10l.paragraphs[0].font.color.rgb = ACCENT_MID
tf10l.paragraphs[0].font.name = '微软雅黑'
add_para(tf10l, "  ② 现代转化缺乏社会学、心理学的量化实证支撑", font_size=15, color=ACCENT_MID, space_before=Pt(16))

# 右侧：未来方向
card_r = add_shape(slide10, Inches(6.8), Inches(2.2), Inches(5.8), Inches(4.3), CARD_BG, LINE_COLOR, 1.5)
add_textbox(slide10, Inches(7.1), Inches(2.35), Inches(5), Inches(0.4),
            "▎ 未来方向", font_size=18, bold=True, color=ACCENT_GOLD)

tb10r = add_textbox(slide10, Inches(7.1), Inches(2.95), Inches(5.3), Inches(3.3), "", font_size=15, color=ACCENT_MID)
tf10r = tb10r.text_frame
tf10r.word_wrap = True
tf10r.paragraphs[0].text = "  ◆  历史人类学田野调查 + 数字人文（NLP/GIS）"
tf10r.paragraphs[0].font.size = Pt(15)
tf10r.paragraphs[0].font.color.rgb = ACCENT_MID
tf10r.paragraphs[0].font.name = '微软雅黑'
add_para(tf10r, "  ◆  与AI伦理、基因编辑等前沿科技伦理跨界对话", font_size=15, color=ACCENT_MID, space_before=Pt(14))
add_para(tf10r, "  ◆  比较宗教学视角下的「全球伦理」范式建构", font_size=15, color=ACCENT_MID, space_before=Pt(14))

add_bottom_bar(slide10)
add_slide_number(slide10, 10)

# ============================================================
# 第11页：致谢
# ============================================================
slide11 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide11, BG_COLOR)
add_top_bar(slide11)

# 中心装饰框
center_card = add_shape(slide11, Inches(3.0), Inches(1.8), Inches(7.3), Inches(4.0), CARD_BG, ACCENT_GOLD, 2.5)

# 致谢
add_textbox(slide11, Inches(3.0), Inches(2.3), Inches(7.3), Inches(1.0),
            "致谢", font_size=40, bold=True, color=ACCENT_GOLD, alignment=PP_ALIGN.CENTER)

# 分割线
add_decorative_line(slide11, Inches(5.5), Inches(3.4), Inches(2.3), ACCENT_GOLD, Pt(2))

add_textbox(slide11, Inches(3.0), Inches(3.7), Inches(7.3), Inches(0.8),
            "恳请各位评委老师批评指正！",
            font_size=22, bold=True, color=ACCENT_DARK, alignment=PP_ALIGN.CENTER)

add_textbox(slide11, Inches(3.0), Inches(4.6), Inches(7.3), Inches(0.6),
            "释果顺  ·  2026年6月",
            font_size=16, color=ACCENT_MID, alignment=PP_ALIGN.CENTER)

# 底部装饰
add_rect(slide11, Inches(5.5), Inches(6.5), Inches(2.3), Inches(0.04), ACCENT_GOLD)

add_slide_number(slide11, 11)

# ========== 保存 ==========
output_path = "/home/admin/.openclaw/workspace/毕业论文答辩_释果顺_重新排版.pptx"
prs.save(output_path)
print(f"✅ PPT已保存到: {output_path}")
