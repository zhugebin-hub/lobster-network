#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
探案风格教学 PPT 模板生成器
"""

from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# 创建演示文稿
prs = Presentation()
prs.slide_width = Cm(33.867)
prs.slide_height = Cm(19.05)

# 探案风格颜色
COLORS = {
    'bg_dark': (20, 20, 30),          # 深蓝色背景
    'bg_light': (35, 35, 50),         # 浅蓝灰背景
    'accent_gold': (212, 175, 55),    # 金色（线索/重点）
    'accent_red': (180, 50, 50),      # 红色（重要/警告）
    'accent_green': (80, 180, 100),   # 绿色（正确/解决）
    'text_light': (230, 230, 230),    # 浅色文字
    'text_dim': (150, 150, 160),      # 暗淡文字
    'white': (255, 255, 255),         # 白色
    'black': (10, 10, 15),            # 黑色
}

def set_background(slide, color=COLORS['bg_dark']):
    """设置背景"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*color)

def add_magnifying_glass(slide, x, y, size=Cm(3)):
    """添加放大镜装饰"""
    # 镜片
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, x, y, size, size
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*COLORS['bg_light'])
    shape.fill.transparency = 0.3
    shape.line.color.rgb = RGBColor(*COLORS['accent_gold'])
    shape.line.width = Pt(4)
    
    # 镜柄
    handle_x = x + size * 0.7
    handle_y = y + size * 0.7
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, handle_x, handle_y, size * 0.3, size * 0.8
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*COLORS['accent_gold'])
    shape.line.fill.background()
    shape.rotation = 45

def add_clue_marker(slide, x, y, text=""):
    """添加线索标记"""
    # 圆形标记
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, x, y, Cm(1.5), Cm(1.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*COLORS['accent_gold'])
    shape.line.color.rgb = RGBColor(*COLORS['white'])
    shape.line.width = Pt(2)
    
    # 数字
    if text:
        textbox = slide.shapes.add_textbox(x + Cm(0.3), y + Cm(0.3), Cm(0.9), Cm(0.9))
        tf = textbox.text_frame
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = RGBColor(*COLORS['bg_dark'])
        p.font.name = 'Microsoft YaHei'
        p.alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

def add_header(slide, title, subtitle=""):
    """添加页眉"""
    # 顶部装饰线
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Cm(0), Cm(0), prs.slide_width, Cm(0.3)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*COLORS['accent_gold'])
    shape.line.fill.background()
    
    # 标题
    textbox = slide.shapes.add_textbox(Cm(1.5), Cm(0.5), Cm(25), Cm(1))
    tf = textbox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = RGBColor(*COLORS['accent_gold'])
    p.font.name = 'Microsoft YaHei'
    
    if subtitle:
        textbox = slide.shapes.add_textbox(Cm(22), Cm(0.6), Cm(10), Cm(0.8))
        tf = textbox.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(*COLORS['text_dim'])
        p.font.name = 'Microsoft YaHei'
        p.alignment = PP_ALIGN.RIGHT

def add_footer(slide, page_num=""):
    """添加页脚"""
    # 底线
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Cm(2), prs.slide_height - Cm(1.2), prs.slide_width - Cm(4), Cm(0.05)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*COLORS['accent_gold'])
    shape.line.fill.background()
    
    # 页码
    textbox = slide.shapes.add_textbox(prs.slide_width - Cm(3), prs.slide_height - Cm(1.5), Cm(2), Cm(0.5))
    tf = textbox.text_frame
    p = tf.paragraphs[0]
    p.text = page_num
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(*COLORS['text_dim'])
    p.font.name = 'Microsoft YaHei'
    p.alignment = PP_ALIGN.RIGHT
    
    # 装饰图标（放大镜）
    add_magnifying_glass(slide, Cm(2.5), prs.slide_height - Cm(1.8), Cm(0.8))

def add_title_slide(prs, title, subtitle, footer=""):
    """封面页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, COLORS['bg_dark'])
    
    # 背景装饰 - 网格线
    for i in range(10):
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Cm(i * 3.5), Cm(0), Cm(0.05), prs.slide_height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*COLORS['bg_light'])
        shape.fill.transparency = 0.8
        shape.line.fill.background()
    
    # 中央标题框
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Cm(5), Cm(5), prs.slide_width - Cm(10), Cm(6)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*COLORS['bg_light'])
    shape.fill.transparency = 0.5
    shape.line.color.rgb = RGBColor(*COLORS['accent_gold'])
    shape.line.width = Pt(3)
    
    # 主标题
    textbox = slide.shapes.add_textbox(Cm(6), Cm(6), prs.slide_width - Cm(12), Cm(2.5))
    tf = textbox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(50)
    p.font.bold = True
    p.font.color.rgb = RGBColor(*COLORS['accent_gold'])
    p.font.name = 'Microsoft YaHei'
    p.alignment = PP_ALIGN.CENTER
    
    # 副标题
    textbox = slide.shapes.add_textbox(Cm(6), Cm(8.5), prs.slide_width - Cm(12), Cm(1.5))
    tf = textbox.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(*COLORS['text_light'])
    p.font.name = 'Microsoft YaHei'
    p.alignment = PP_ALIGN.CENTER
    
    # 底部信息
    if footer:
        textbox = slide.shapes.add_textbox(Cm(6), prs.slide_height - Cm(3), prs.slide_width - Cm(12), Cm(1))
        tf = textbox.text_frame
        p = tf.paragraphs[0]
        p.text = footer
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(*COLORS['text_dim'])
        p.font.name = 'Microsoft YaHei'
        p.alignment = PP_ALIGN.CENTER
    
    # 装饰放大镜
    add_magnifying_glass(slide, Cm(25), Cm(12), Cm(4))
    
    return slide

def add_case_brief_slide(prs, title, case_info, page_num=""):
    """案情简报页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, COLORS['bg_light'])
    add_header(slide, title, "案情简报")
    
    # 案情框
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Cm(2), Cm(2.5), prs.slide_width - Cm(4), Cm(13)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*COLORS['bg_dark'])
    shape.line.color.rgb = RGBColor(*COLORS['accent_gold'])
    shape.line.width = Pt(2)
    shape.line.dash_style = 4  # 虚线
    
    # 案情内容
    y = Cm(3.5)
    for item in case_info:
        if isinstance(item, dict):
            if item.get('type') == 'label':
                textbox = slide.shapes.add_textbox(Cm(3), y, Cm(3), Cm(0.8))
                tf = textbox.text_frame
                p = tf.paragraphs[0]
                p.text = item['text']
                p.font.size = Pt(18)
                p.font.bold = True
                p.font.color.rgb = RGBColor(*COLORS['accent_gold'])
                p.font.name = 'Microsoft YaHei'
            elif item.get('type') == 'content':
                textbox = slide.shapes.add_textbox(Cm(6), y, Cm(24), Cm(0.8))
                tf = textbox.text_frame
                p = tf.paragraphs[0]
                p.text = item['text']
                p.font.size = Pt(18)
                p.font.color.rgb = RGBColor(*COLORS['text_light'])
                p.font.name = 'Microsoft YaHei'
                y += Cm(1)
        else:
            y += Cm(0.5)
    
    add_footer(slide, page_num)
    return slide

def add_clue_slide(prs, title, clues, page_num=""):
    """线索分析页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, COLORS['bg_light'])
    add_header(slide, title, "线索分析")
    
    # 线索卡片
    clue_width = Cm(13)
    clue_height = Cm(5)
    start_x = Cm(2.5)
    start_y = Cm(3)
    
    for i, clue in enumerate(clues):
        col = i % 2
        row = i // 2
        x = start_x + col * (clue_width + Cm(1))
        y = start_y + row * (clue_height + Cm(0.5))
        
        # 卡片背景
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, y, clue_width, clue_height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*COLORS['bg_dark'])
        shape.line.color.rgb = RGBColor(*COLORS['accent_gold'])
        shape.line.width = Pt(2)
        
        # 线索编号
        add_clue_marker(slide, x + Cm(0.3), y + Cm(0.3), str(i + 1))
        
        # 线索标题
        textbox = slide.shapes.add_textbox(x + Cm(1.5), y + Cm(0.5), clue_width - Cm(2), Cm(0.8))
        tf = textbox.text_frame
        p = tf.paragraphs[0]
        p.text = clue['title']
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = RGBColor(*COLORS['accent_gold'])
        p.font.name = 'Microsoft YaHei'
        
        # 线索内容
        textbox = slide.shapes.add_textbox(x + Cm(1.5), y + Cm(1.5), clue_width - Cm(2), clue_height - Cm(2))
        tf = textbox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = clue['content']
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(*COLORS['text_light'])
        p.font.name = 'Microsoft YaHei'
    
    add_footer(slide, page_num)
    return slide

def add_suspect_slide(prs, title, suspects, page_num=""):
    """嫌疑人分析页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, COLORS['bg_light'])
    add_header(slide, title, "嫌疑人分析")
    
    # 嫌疑人卡片
    suspect_width = Cm(9)
    suspect_height = Cm(7)
    start_x = Cm(3)
    start_y = Cm(3)
    
    for i, suspect in enumerate(suspects):
        x = start_x + i * (suspect_width + Cm(1))
        
        # 卡片背景
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, x, start_y, suspect_width, suspect_height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*COLORS['bg_dark'])
        shape.line.color.rgb = RGBColor(*COLORS['accent_gold'])
        shape.line.width = Pt(2)
        
        # 头像占位
        avatar_size = Cm(3)
        avatar_x = x + suspect_width / 2 - avatar_size / 2
        shape = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, avatar_x, start_y + Cm(0.5), avatar_size, avatar_size
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*COLORS['bg_light'])
        shape.line.color.rgb = RGBColor(*COLORS['text_dim'])
        shape.line.width = Pt(2)
        shape.line.dash_style = 4
        
        # 姓名
        textbox = slide.shapes.add_textbox(x + Cm(0.5), start_y + Cm(4), suspect_width - Cm(1), Cm(0.8))
        tf = textbox.text_frame
        p = tf.paragraphs[0]
        p.text = suspect['name']
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = RGBColor(*COLORS['accent_gold'])
        p.font.name = 'Microsoft YaHei'
        p.alignment = PP_ALIGN.CENTER
        
        # 信息
        textbox = slide.shapes.add_textbox(x + Cm(0.5), start_y + Cm(4.8), suspect_width - Cm(1), Cm(2))
        tf = textbox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = suspect['info']
        p.font.size = Pt(13)
        p.font.color.rgb = RGBColor(*COLORS['text_light'])
        p.font.name = 'Microsoft YaHei'
    
    add_footer(slide, page_num)
    return slide

def add_content_slide(prs, title, content_items, page_num=""):
    """通用内容页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, COLORS['bg_light'])
    add_header(slide, title)
    
    y = Cm(3)
    for item in content_items:
        if isinstance(item, dict):
            if item.get('type') == 'highlight':
                # 高亮框
                shape = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE, Cm(2), y, prs.slide_width - Cm(4), Cm(1.5)
                )
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor(*COLORS['bg_dark'])
                shape.line.color.rgb = RGBColor(*COLORS['accent_gold'])
                shape.line.width = Pt(2)
                
                textbox = slide.shapes.add_textbox(Cm(3), y + Cm(0.3), prs.slide_width - Cm(6), Cm(0.9))
                tf = textbox.text_frame
                p = tf.paragraphs[0]
                p.text = item['text']
                p.font.size = Pt(20)
                p.font.bold = True
                p.font.color.rgb = RGBColor(*COLORS['accent_gold'])
                p.font.name = 'Microsoft YaHei'
                y += Cm(2)
            elif item.get('type') == 'bullet':
                textbox = slide.shapes.add_textbox(Cm(3), y, prs.slide_width - Cm(6), Cm(0.8))
                tf = textbox.text_frame
                p = tf.paragraphs[0]
                p.text = "🔍 " + item['text']
                p.font.size = Pt(18)
                p.font.color.rgb = RGBColor(*COLORS['text_light'])
                p.font.name = 'Microsoft YaHei'
                y += Cm(1)
            elif item.get('type') == 'conclusion':
                # 结论框
                shape = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE, Cm(2), y, prs.slide_width - Cm(4), Cm(2)
                )
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor(*COLORS['accent_green'])
                shape.fill.transparency = 0.8
                shape.line.color.rgb = RGBColor(*COLORS['accent_green'])
                shape.line.width = Pt(2)
                
                textbox = slide.shapes.add_textbox(Cm(3), y + Cm(0.5), prs.slide_width - Cm(6), Cm(1))
                tf = textbox.text_frame
                p = tf.paragraphs[0]
                p.text = "✓ " + item['text']
                p.font.size = Pt(20)
                p.font.bold = True
                p.font.color.rgb = RGBColor(*COLORS['accent_green'])
                p.font.name = 'Microsoft YaHei'
                y += Cm(2.5)
        else:
            textbox = slide.shapes.add_textbox(Cm(3), y, prs.slide_width - Cm(6), Cm(0.8))
            tf = textbox.text_frame
            p = tf.paragraphs[0]
            p.text = item
            p.font.size = Pt(18)
            p.font.color.rgb = RGBColor(*COLORS['text_light'])
            p.font.name = 'Microsoft YaHei'
            y += Cm(1)
    
    add_footer(slide, page_num)
    return slide

def add_mystery_slide(prs, title, question, page_num=""):
    """悬念提问页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, COLORS['bg_dark'])
    
    # 问号装饰
    textbox = slide.shapes.add_textbox(Cm(12), Cm(3), Cm(10), Cm(10))
    tf = textbox.text_frame
    p = tf.paragraphs[0]
    p.text = "?"
    p.font.size = Pt(200)
    p.font.color.rgb = RGBColor(*COLORS['accent_gold'])
    p.font.name = 'Arial'
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_after = Pt(0)
    
    # 问题框
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Cm(5), Cm(10), prs.slide_width - Cm(10), Cm(4)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*COLORS['bg_light'])
    shape.line.color.rgb = RGBColor(*COLORS['accent_gold'])
    shape.line.width = Pt(3)
    
    # 问题文字
    textbox = slide.shapes.add_textbox(Cm(6), Cm(11), prs.slide_width - Cm(12), Cm(2.5))
    tf = textbox.text_frame
    p = tf.paragraphs[0]
    p.text = question
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(*COLORS['accent_gold'])
    p.font.name = 'Microsoft YaHei'
    p.alignment = PP_ALIGN.CENTER
    
    add_footer(slide, page_num)
    return slide

def add_solution_slide(prs, title, solution, page_num=""):
    """真相揭示页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, COLORS['bg_light'])
    add_header(slide, title, "真相大白")
    
    # 解决方案框
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Cm(3), Cm(3), prs.slide_width - Cm(6), Cm(12)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*COLORS['bg_dark'])
    shape.line.color.rgb = RGBColor(*COLORS['accent_green'])
    shape.line.width = Pt(3)
    
    # 成功标记
    textbox = slide.shapes.add_textbox(Cm(4), Cm(4), Cm(2), Cm(2))
    tf = textbox.text_frame
    p = tf.paragraphs[0]
    p.text = "✓"
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = RGBColor(*COLORS['accent_green'])
    p.font.name = 'Arial'
    
    # 内容
    y = Cm(4.5)
    for item in solution:
        textbox = slide.shapes.add_textbox(Cm(7), y, prs.slide_width - Cm(11), Cm(1))
        tf = textbox.text_frame
        p = tf.paragraphs[0]
        p.text = item
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(*COLORS['text_light'])
        p.font.name = 'Microsoft YaHei'
        y += Cm(1.2)
    
    add_footer(slide, page_num)
    return slide

def add_end_slide(prs, title, subtitle=""):
    """结束页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, COLORS['bg_dark'])
    
    # 中央文字
    textbox = slide.shapes.add_textbox(Cm(6), Cm(7), prs.slide_width - Cm(12), Cm(3))
    tf = textbox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(50)
    p.font.bold = True
    p.font.color.rgb = RGBColor(*COLORS['accent_gold'])
    p.font.name = 'Microsoft YaHei'
    p.alignment = PP_ALIGN.CENTER
    
    if subtitle:
        textbox = slide.shapes.add_textbox(Cm(6), Cm(10), prs.slide_width - Cm(12), Cm(2))
        tf = textbox.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor(*COLORS['text_dim'])
        p.font.name = 'Microsoft YaHei'
        p.alignment = PP_ALIGN.CENTER
    
    # 装饰
    add_magnifying_glass(slide, Cm(25), Cm(12), Cm(4))
    
    return slide

# ============ 生成模板幻灯片 ============

# 1. 封面
add_title_slide(prs,
    title="案件调查",
    subtitle="探案风格教学 PPT 模板",
    footer="适用于各类课程的情境化教学"
)

# 2. 学习目标
add_content_slide(prs, "本节任务", [
    {"type": "highlight", "text": "你需要完成以下调查任务："},
    "",
    {"type": "bullet", "text": "阅读案情简报，了解基本事实"},
    {"type": "bullet", "text": "分析现场线索，找出关键证据"},
    {"type": "bullet", "text": "审问嫌疑人，发现矛盾之处"},
    {"type": "bullet", "text": "推理真相，找出真凶"},
    {"type": "bullet", "text": "总结案件，形成完整报告"},
], "1/10")

# 3. 案情简报
add_case_brief_slide(prs, "案情简介", [
    {"type": "label", "text": "案件名称："},
    {"type": "content", "text": "点击编辑案件名称"},
    "",
    {"type": "label", "text": "发生时间："},
    {"type": "content", "text": "点击编辑时间"},
    "",
    {"type": "label", "text": "发生地点："},
    {"type": "content", "text": "点击编辑地点"},
    "",
    {"type": "label", "text": "案件概述："},
    {"type": "content", "text": "点击编辑案件详细描述..."},
], "2/10")

# 4. 线索分析
add_clue_slide(prs, "现场线索", [
    {"title": "线索一", "content": "点击编辑线索内容..."},
    {"title": "线索二", "content": "点击编辑线索内容..."},
    {"title": "线索三", "content": "点击编辑线索内容..."},
    {"title": "线索四", "content": "点击编辑线索内容..."},
], "3/10")

# 5. 嫌疑人
add_suspect_slide(prs, "嫌疑人分析", [
    {"name": "嫌疑人 A", "info": "点击编辑\n年龄/职业/动机..."},
    {"name": "嫌疑人 B", "info": "点击编辑\n年龄/职业/动机..."},
    {"name": "嫌疑人 C", "info": "点击编辑\n年龄/职业/动机..."},
], "4/10")

# 6. 悬念提问
add_mystery_slide(prs, "关键问题", "谁是真正的凶手？", "5/10")

# 7. 推理分析
add_content_slide(prs, "推理过程", [
    {"type": "highlight", "text": "根据现有线索，我们来逐步推理："},
    "",
    {"type": "bullet", "text": "首先，从时间线上分析..."},
    {"type": "bullet", "text": "其次，从动机角度考虑..."},
    {"type": "bullet", "text": "再次，从证据链验证..."},
    {"type": "bullet", "text": "最后，排除不可能的选项..."},
    "",
    {"type": "conclusion", "text": "结论：点击编辑你的推理结论"},
], "6/10")

# 8. 真相揭示
add_solution_slide(prs, "真相大白", [
    "真凶是：点击编辑",
    "",
    "作案动机：点击编辑",
    "",
    "作案手法：点击编辑",
    "",
    "关键证据：点击编辑",
], "7/10")

# 9. 知识总结
add_content_slide(prs, "案件总结", [
    {"type": "highlight", "text": "通过本案，我们学习了："},
    "",
    {"type": "bullet", "text": "知识点一：点击编辑"},
    {"type": "bullet", "text": "知识点二：点击编辑"},
    {"type": "bullet", "text": "知识点三：点击编辑"},
    {"type": "bullet", "text": "技能点：点击编辑"},
    "",
    {"type": "conclusion", "text": "核心素养：逻辑推理能力、证据分析能力"},
], "8/10")

# 10. 课后任务
add_content_slide(prs, "课后拓展", [
    {"type": "highlight", "text": "课后任务："},
    "",
    {"type": "bullet", "text": "基础题：整理本案的完整证据链"},
    {"type": "bullet", "text": "提高题：设计一个类似的推理案例"},
    {"type": "bullet", "text": "拓展题：阅读推荐书籍/观看相关影片"},
    "",
    {"type": "bullet", "text": "下节预告：新的案件等待你来破解..."},
], "9/10")

# 11. 结束页
add_end_slide(prs, "调查结束", "感谢参与本次案件调查")

# 保存文件
output_path = '/home/admin/.openclaw/workspace/探案风格教学 PPT 模板.pptx'
prs.save(output_path)
print(f"PPT 模板已生成：{output_path}")
print(f"共 {len(prs.slides)} 页")
