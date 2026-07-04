#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
初识大数据教学 PPT 生成器
"""

from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# 创建演示文稿
prs = Presentation()
prs.slide_width = Cm(33.867)
prs.slide_height = Cm(19.05)

# 颜色定义
COLORS = {
    'primary': (41, 128, 185),      # 蓝色
    'secondary': (52, 152, 219),    # 浅蓝
    'dark': (44, 62, 80),           # 深蓝灰
    'white': (255, 255, 255),       # 白色
    'accent': (231, 76, 60),        # 红色
    'success': (39, 174, 96),       # 绿色
    'warning': (243, 156, 18),      # 橙色
    'light': (236, 240, 241),       # 浅灰
    'gray': (127, 140, 141),        # 灰色
}

def set_background(slide, color=COLORS['white']):
    """设置背景"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*color)

def add_header(slide, title, subtitle=""):
    """添加页眉"""
    # 顶部色带
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Cm(0), Cm(0), prs.slide_width, Cm(1.8)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*COLORS['primary'])
    shape.line.fill.background()
    
    # 标题
    textbox = slide.shapes.add_textbox(Cm(1.5), Cm(0.4), Cm(20), Cm(1))
    tf = textbox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(*COLORS['white'])
    p.font.name = 'Microsoft YaHei'
    
    if subtitle:
        textbox = slide.shapes.add_textbox(Cm(22), Cm(0.5), Cm(10), Cm(0.8))
        tf = textbox.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(*COLORS['white'])
        p.font.name = 'Microsoft YaHei'
        p.alignment = PP_ALIGN.RIGHT

def add_footer(slide, page_num=""):
    """添加页脚"""
    # 底线
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Cm(2), prs.slide_height - Cm(1), prs.slide_width - Cm(4), Cm(0.05)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*COLORS['secondary'])
    shape.line.fill.background()
    
    # 页码
    textbox = slide.shapes.add_textbox(prs.slide_width - Cm(3), prs.slide_height - Cm(1.3), Cm(2), Cm(0.5))
    tf = textbox.text_frame
    p = tf.paragraphs[0]
    p.text = page_num
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(*COLORS['gray'])
    p.font.name = 'Microsoft YaHei'
    p.alignment = PP_ALIGN.RIGHT

def add_title_slide(prs, title, subtitle, footer=""):
    """封面页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    
    # 背景装饰
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Cm(0), Cm(0), prs.slide_width, Cm(5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*COLORS['primary'])
    shape.line.fill.background()
    
    # 主标题
    textbox = slide.shapes.add_textbox(Cm(2), Cm(1.5), Cm(28), Cm(3))
    tf = textbox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = RGBColor(*COLORS['white'])
    p.font.name = 'Microsoft YaHei'
    
    # 副标题
    textbox = slide.shapes.add_textbox(Cm(2), Cm(5), Cm(20), Cm(1.5))
    tf = textbox.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(*COLORS['white'])
    p.font.name = 'Microsoft YaHei'
    
    # 底部信息
    if footer:
        textbox = slide.shapes.add_textbox(Cm(2), prs.slide_height - Cm(3), Cm(15), Cm(1))
        tf = textbox.text_frame
        p = tf.paragraphs[0]
        p.text = footer
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(*COLORS['gray'])
        p.font.name = 'Microsoft YaHei'
    
    return slide

def add_content_slide(prs, title, content_items, page_num=""):
    """内容页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_header(slide, title)
    
    # 内容
    y = Cm(3)
    for item in content_items:
        if isinstance(item, dict):
            if item.get('type') == 'title':
                textbox = slide.shapes.add_textbox(Cm(2), y, Cm(30), Cm(1))
                tf = textbox.text_frame
                p = tf.paragraphs[0]
                p.text = item['text']
                p.font.size = Pt(24)
                p.font.bold = True
                p.font.color.rgb = RGBColor(*COLORS['primary'])
                p.font.name = 'Microsoft YaHei'
                y += Cm(1.2)
            elif item.get('type') == 'bullet':
                textbox = slide.shapes.add_textbox(Cm(2.5), y, Cm(29), Cm(0.8))
                tf = textbox.text_frame
                p = tf.paragraphs[0]
                p.text = "• " + item['text']
                p.font.size = Pt(20)
                p.font.color.rgb = RGBColor(*COLORS['dark'])
                p.font.name = 'Microsoft YaHei'
                y += Cm(1)
            elif item.get('type') == 'highlight':
                # 高亮框
                shape = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE, Cm(2), y, Cm(30), Cm(1.5)
                )
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor(*COLORS['light'])
                shape.line.color.rgb = RGBColor(*COLORS['primary'])
                shape.line.width = Pt(2)
                
                textbox = slide.shapes.add_textbox(Cm(2.5), y + Cm(0.3), Cm(29), Cm(0.9))
                tf = textbox.text_frame
                p = tf.paragraphs[0]
                p.text = item['text']
                p.font.size = Pt(20)
                p.font.bold = True
                p.font.color.rgb = RGBColor(*COLORS['primary'])
                p.font.name = 'Microsoft YaHei'
                y += Cm(2)
            elif item.get('type') == 'box':
                # 信息框
                shape = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE, Cm(2), y, Cm(30), item.get('height', Cm(2))
                )
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor(*COLORS.get(item.get('color', 'light'), COLORS['light']))
                shape.line.fill.background()
                
                textbox = slide.shapes.add_textbox(Cm(2.5), y + Cm(0.3), Cm(29), Cm(1.5))
                tf = textbox.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = item['text']
                p.font.size = Pt(18)
                p.font.color.rgb = RGBColor(*COLORS['dark'])
                p.font.name = 'Microsoft YaHei'
                y += item.get('height', Cm(2)) + Cm(0.3)
        else:
            textbox = slide.shapes.add_textbox(Cm(2), y, Cm(30), Cm(0.8))
            tf = textbox.text_frame
            p = tf.paragraphs[0]
            p.text = item
            p.font.size = Pt(20)
            p.font.color.rgb = RGBColor(*COLORS['dark'])
            p.font.name = 'Microsoft YaHei'
            y += Cm(1)
    
    add_footer(slide, page_num)
    return slide

def add_comparison_slide(prs, title, items, page_num=""):
    """对比页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_header(slide, title)
    
    # 创建对比表格
    col_width = Cm(14)
    start_x = Cm(2)
    start_y = Cm(3)
    row_height = Cm(2)
    
    # 表头
    for i, header in enumerate(items['headers']):
        x = start_x + i * (col_width + Cm(0.5))
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, x, start_y, col_width, Cm(1.2)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*COLORS['primary'])
        shape.line.fill.background()
        
        textbox = slide.shapes.add_textbox(x + Cm(0.5), start_y + Cm(0.2), col_width - Cm(1), Cm(0.8))
        tf = textbox.text_frame
        p = tf.paragraphs[0]
        p.text = header
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = RGBColor(*COLORS['white'])
        p.font.name = 'Microsoft YaHei'
        p.alignment = PP_ALIGN.CENTER
    
    # 内容行
    for row_idx, row_data in enumerate(items['rows']):
        y = start_y + Cm(1.5) + row_idx * (row_height + Cm(0.3))
        for col_idx, cell in enumerate(row_data):
            x = start_x + col_idx * (col_width + Cm(0.5))
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, x, y, col_width, row_height
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(*COLORS['light'])
            shape.line.color.rgb = RGBColor(*COLORS['gray'])
            shape.line.width = Pt(1)
            
            textbox = slide.shapes.add_textbox(x + Cm(0.5), y + Cm(0.3), col_width - Cm(1), row_height - Cm(0.6))
            tf = textbox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = cell
            p.font.size = Pt(16)
            p.font.color.rgb = RGBColor(*COLORS['dark'])
            p.font.name = 'Microsoft YaHei'
    
    add_footer(slide, page_num)
    return slide

def add_diagram_slide(prs, title, diagram_type, page_num=""):
    """图示页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_header(slide, title)
    
    if diagram_type == '4V':
        # 4V 特征图
        positions = [
            (Cm(4), Cm(4), "Volume\n大量化", COLORS['primary']),
            (Cm(18), Cm(4), "Variety\n多样化", COLORS['secondary']),
            (Cm(4), Cm(10), "Velocity\n快速化", COLORS['success']),
            (Cm(18), Cm(10), "Value\n价值化", COLORS['warning']),
        ]
        
        for x, y, text, color in positions:
            shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Cm(12), Cm(5)
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(*color)
            shape.line.fill.background()
            
            textbox = slide.shapes.add_textbox(x + Cm(1), y + Cm(1), Cm(10), Cm(3))
            tf = textbox.text_frame
            p = tf.paragraphs[0]
            p.text = text
            p.font.size = Pt(24)
            p.font.bold = True
            p.font.color.rgb = RGBColor(*COLORS['white'])
            p.font.name = 'Microsoft YaHei'
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    add_footer(slide, page_num)
    return slide

# ============ 生成幻灯片 ============

# 1. 封面
add_title_slide(prs, 
    title="初识大数据",
    subtitle="Big Data 基础入门",
    footer="中职/高职计算机专业 • 2026 年 4 月"
)

# 2. 学习目标
add_content_slide(prs, "学习目标", [
    {"type": "bullet", "text": "了解大数据的定义和发展背景"},
    {"type": "bullet", "text": "掌握大数据的 4V 特征"},
    {"type": "bullet", "text": "认识大数据的典型应用场景"},
    {"type": "bullet", "text": "了解大数据处理的基本流程"},
    {"type": "bullet", "text": "认识常见的大数据技术工具"},
], "1/12")

# 3. 什么是大数据
add_content_slide(prs, "什么是大数据？", [
    {"type": "highlight", "text": "大数据 (Big Data) 是指无法在一定时间范围内用常规软件工具进行捕捉、管理和处理的数据集合。"},
    "",
    {"type": "title", "text": "简单理解："},
    {"type": "bullet", "text": "数据量太大，传统数据库存不下"},
    {"type": "bullet", "text": "数据类型太多，传统方法处理不了"},
    {"type": "bullet", "text": "数据产生太快，传统技术跟不上"},
    "",
    {"type": "box", "text": "大数据不是单纯的技术问题，而是需要新的处理模式才能具有更强的决策力、洞察发现力和流程优化能力。", "color": 'light', "height": Cm(2)},
], "2/12")

# 4. 大数据的发展背景
add_content_slide(prs, "大数据的发展背景", [
    {"type": "title", "text": "为什么大数据会兴起？"},
    "",
    {"type": "bullet", "text": "互联网普及：全球网民超过 50 亿，每天产生海量数据"},
    {"type": "bullet", "text": "移动设备：智能手机、平板电脑随时产生数据"},
    {"type": "bullet", "text": "物联网：传感器、智能设备 24 小时采集数据"},
    {"type": "bullet", "text": "社交媒体：微信、微博、抖音等内容爆发"},
    {"type": "bullet", "text": "存储成本下降：硬盘价格越来越便宜"},
    {"type": "bullet", "text": "计算能力提升：分布式计算技术成熟"},
], "3/12")

# 5. 大数据的 4V 特征
add_diagram_slide(prs, "大数据的 4V 特征", "4V", "4/12")

# 6. 4V 特征详解
add_content_slide(prs, "4V 特征详解", [
    {"type": "title", "text": "Volume (大量化)"},
    {"type": "bullet", "text": "数据量从 TB 级跃升到 PB、EB 级"},
    {"type": "bullet", "text": "例：YouTube 每分钟上传 500 小时视频"},
    "",
    {"type": "title", "text": "Variety (多样化)"},
    {"type": "bullet", "text": "结构化数据 (数据库表格)"},
    {"type": "bullet", "text": "半结构化数据 (XML、JSON)"},
    {"type": "bullet", "text": "非结构化数据 (文本、图片、音频、视频)"},
    "",
    {"type": "title", "text": "Velocity (快速化)"},
    {"type": "bullet", "text": "数据产生和更新速度快"},
    {"type": "bullet", "text": "需要实时或近实时处理"},
    "",
    {"type": "title", "text": "Value (价值化)"},
    {"type": "bullet", "text": "数据价值密度低，需要挖掘"},
    {"type": "bullet", "text": "例：1 小时监控视频，有价值的可能只有几秒"},
], "5/12")

# 7. 大数据 vs 传统数据
add_comparison_slide(prs, "大数据 vs 传统数据", {
    'headers': ['对比维度', '传统数据', '大数据'],
    'rows': [
        ['数据规模', 'GB-TB 级', 'PB-EB 级'],
        ['数据类型', '结构化为主', '结构化 + 非结构化'],
        ['处理方式', '集中式处理', '分布式处理'],
        ['响应时间', '分钟 - 小时级', '秒级 - 实时'],
        ['存储成本', '较高', '较低 (廉价硬件)'],
        ['应用场景', '业务系统', '数据分析、智能决策'],
    ]
}, "6/12")

# 8. 大数据应用场景
add_content_slide(prs, "大数据应用场景", [
    {"type": "title", "text": "电商领域"},
    {"type": "bullet", "text": "个性化推荐 (淘宝猜你喜欢)"},
    {"type": "bullet", "text": "精准营销 (优惠券推送)"},
    {"type": "bullet", "text": "库存预测 (销量预测)"},
    "",
    {"type": "title", "text": "金融领域"},
    {"type": "bullet", "text": "风险控制 (信用评估)"},
    {"type": "bullet", "text": "欺诈检测 (异常交易识别)"},
    {"type": "bullet", "text": "智能投顾 (投资建议)"},
    "",
    {"type": "title", "text": "医疗领域"},
    {"type": "bullet", "text": "疾病预测 (流行病学分析)"},
    {"type": "bullet", "text": "辅助诊断 (医学影像分析)"},
    {"type": "bullet", "text": "药物研发 (基因数据分析)"},
], "7/12")

# 9. 大数据应用场景 (续)
add_content_slide(prs, "大数据应用场景 (续)", [
    {"type": "title", "text": "交通领域"},
    {"type": "bullet", "text": "智能导航 (实时路况分析)"},
    {"type": "bullet", "text": "出行预测 (滴滴/Uber 调度)"},
    {"type": "bullet", "text": "自动驾驶 (传感器数据处理)"},
    "",
    {"type": "title", "text": "教育领域"},
    {"type": "bullet", "text": "个性化学习 (学习行为分析)"},
    {"type": "bullet", "text": "教学质量评估 (成绩分析)"},
    {"type": "bullet", "text": "辍学预警 (学生行为监测)"},
    "",
    {"type": "title", "text": "政务领域"},
    {"type": "bullet", "text": "智慧城市 (城市运行监测)"},
    {"type": "bullet", "text": "公共安全 (犯罪预测)"},
    {"type": "bullet", "text": "应急管理 (灾害预警)"},
], "8/12")

# 10. 大数据处理流程
add_content_slide(prs, "大数据处理流程", [
    {"type": "highlight", "text": "数据采集 → 数据存储 → 数据处理 → 数据分析 → 数据可视化"},
    "",
    {"type": "title", "text": "1. 数据采集"},
    {"type": "bullet", "text": "日志采集、传感器数据、API 接口、网络爬虫"},
    "",
    {"type": "title", "text": "2. 数据存储"},
    {"type": "bullet", "text": "HDFS、HBase、MongoDB、云存储"},
    "",
    {"type": "title", "text": "3. 数据处理"},
    {"type": "bullet", "text": "MapReduce、Spark、Flink 分布式计算"},
    "",
    {"type": "title", "text": "4. 数据分析"},
    {"type": "bullet", "text": "统计分析、机器学习、深度学习"},
    "",
    {"type": "title", "text": "5. 数据可视化"},
    {"type": "bullet", "text": "ECharts、Tableau、PowerBI 展示结果"},
], "9/12")

# 11. 常见大数据技术
add_content_slide(prs, "常见大数据技术工具", [
    {"type": "title", "text": "存储技术"},
    {"type": "bullet", "text": "HDFS (Hadoop 分布式文件系统)"},
    {"type": "bullet", "text": "HBase (分布式列式数据库)"},
    {"type": "bullet", "text": "MongoDB (文档数据库)"},
    "",
    {"type": "title", "text": "计算技术"},
    {"type": "bullet", "text": "MapReduce (离线批处理)"},
    {"type": "bullet", "text": "Spark (内存计算，支持批流)"},
    {"type": "bullet", "text": "Flink (实时流处理)"},
    "",
    {"type": "title", "text": "资源调度"},
    {"type": "bullet", "text": "YARN (Hadoop 资源管理器)"},
    {"type": "bullet", "text": "Kubernetes (容器编排)"},
    "",
    {"type": "title", "text": "数据采集"},
    {"type": "bullet", "text": "Flume (日志采集)"},
    {"type": "bullet", "text": "Kafka (消息队列)"},
], "10/12")

# 12. 大数据学习建议
add_content_slide(prs, "大数据学习建议", [
    {"type": "highlight", "text": "学习路线：基础 → 框架 → 实战"},
    "",
    {"type": "title", "text": "基础知识"},
    {"type": "bullet", "text": "Linux 操作系统基础"},
    {"type": "bullet", "text": "Java/Python 编程语言"},
    {"type": "bullet", "text": "数据库基础 (SQL)"},
    {"type": "bullet", "text": "计算机网络基础"},
    "",
    {"type": "title", "text": "核心框架"},
    {"type": "bullet", "text": "Hadoop (HDFS + MapReduce + YARN)"},
    {"type": "bullet", "text": "Spark (RDD + SparkSQL + SparkStreaming)"},
    {"type": "bullet", "text": "Hive (数据仓库)"},
    "",
    {"type": "title", "text": "实践建议"},
    {"type": "bullet", "text": "搭建个人实验环境 (虚拟机/Docker)"},
    {"type": "bullet", "text": "完成开源项目实战"},
    {"type": "bullet", "text": "参加大数据竞赛 (Kaggle、天池)"},
], "11/12")

# 13. 课后练习
add_content_slide(prs, "课后练习与思考", [
    {"type": "title", "text": "基础题"},
    {"type": "bullet", "text": "简述大数据的 4V 特征"},
    {"type": "bullet", "text": "列举 3 个大数据在日常生活中的应用"},
    {"type": "bullet", "text": "大数据与传统数据的主要区别是什么？"},
    "",
    {"type": "title", "text": "提高题"},
    {"type": "bullet", "text": "调研一家你熟悉的公司 (如淘宝、抖音) 是如何使用大数据的"},
    {"type": "bullet", "text": "思考：大数据可能带来哪些隐私和安全问题？"},
    "",
    {"type": "highlight", "text": "拓展阅读：《大数据时代》《数据科学入门》"},
], "12/12")

# 14. 结束页
add_title_slide(prs,
    title="感谢观看",
    subtitle="敬请批评指正",
    footer="初识大数据 • 2026 年 4 月"
)

# 保存文件
output_path = '/home/admin/.openclaw/workspace/初识大数据教学 PPT.pptx'
prs.save(output_path)
print(f"PPT 已生成：{output_path}")
print(f"共 {len(prs.slides)} 页")
