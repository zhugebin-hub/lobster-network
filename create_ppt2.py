#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
根据图片内容创建可编辑的 PPT 文件 - 三课展示促能
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def create_ppt():
    # 创建演示文稿
    prs = Presentation()
    
    # 使用 16:9 的幻灯片
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # ========== 第一页：三课展示促能 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白版式
    
    # 设置背景色为浅灰色
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(245, 245, 245)
    
    # 主标题框
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(4), Inches(0.8))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = "三课展示促能"
    title_para.font.size = Pt(28)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(44, 62, 80)
    
    # 标题左侧红色三角
    triangle = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_TRIANGLE,
        Inches(0.3), Inches(0.6), Inches(0.3), Inches(0.4)
    )
    triangle.fill.solid()
    triangle.fill.fore_color.rgb = RGBColor(231, 76, 60)
    triangle.line.fill.background()
    
    # 标题右侧蓝色横幅
    banner = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(8.5), Inches(0.4), Inches(4.5), Inches(0.8)
    )
    banner.fill.solid()
    banner.fill.fore_color.rgb = RGBColor(41, 128, 185)
    banner.line.fill.background()
    
    # 横幅文字
    banner_text = slide.shapes.add_textbox(Inches(8.7), Inches(0.55), Inches(4), Inches(0.5))
    banner_frame = banner_text.text_frame
    banner_para = banner_frame.paragraphs[0]
    banner_para.text = "管理三课堂"
    banner_para.font.size = Pt(20)
    banner_para.font.bold = True
    banner_para.font.color.rgb = RGBColor(255, 255, 255)
    banner_para.alignment = PP_ALIGN.CENTER
    
    # 顶部红色虚线
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5), Inches(1.15), Inches(12.5), Inches(0.05)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(231, 76, 60)
    line.line.fill.background()
    
    # 定义三个内容板块
    content_boxes = [
        {
            "title": "青蓝共进课",
            "content": "新教师：每周听课量不少于本学科周课时的 2/3，课后及时与师傅交流。\n指导教师：每周至少听徒弟 1 节课，并给予具体评价和改进方向。",
            "x": 1.5,
            "y": 1.8,
            "color": RGBColor(231, 76, 60),
            "num": 1
        },
        {
            "title": "行政视导课",
            "content": "听课要求：全体行政每月至少进入本学科新教师课堂一次，深入一线。\n指导方向：从学科教学、课堂管理和学科研究等方面给予全面专业指导，助力新教师快速适应岗位。",
            "x": 7.5,
            "y": 1.8,
            "color": RGBColor(41, 128, 185),
            "num": 2
        },
        {
            "title": "年级随堂课",
            "content": "听课目的：准确了解各学科、各班级的教学发展现状，把握年级教学脉搏。\n听课要求：年级管理小组全面深入课堂，重点关注学科覆盖率、班级覆盖率和教师覆盖率。\n重点关注：通过课堂观察，多维度了解学生学习情况、学科教学情况和班级管理现状。",
            "x": 4.5,
            "y": 4.2,
            "color": RGBColor(142, 68, 173),
            "num": 3
        }
    ]
    
    # 创建三个内容卡片
    for i, box in enumerate(content_boxes):
        # 卡片背景
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(box["x"]), Inches(box["y"]), Inches(4.5), Inches(2.5)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(255, 255, 255)
        card.line.color.rgb = box["color"]
        card.line.width = Pt(2)
        
        # 卡片标题
        title_box = slide.shapes.add_textbox(
            Inches(box["x"] + 0.3), Inches(box["y"] + 0.3), Inches(3.9), Inches(0.5)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = box["title"]
        title_para.font.size = Pt(16)
        title_para.font.bold = True
        title_para.font.color.rgb = box["color"]
        
        # 卡片内容
        content_box = slide.shapes.add_textbox(
            Inches(box["x"] + 0.3), Inches(box["y"] + 0.8), Inches(3.9), Inches(1.5)
        )
        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        content_para = content_frame.paragraphs[0]
        content_para.text = box["content"]
        content_para.font.size = Pt(11)
        content_para.font.color.rgb = RGBColor(80, 80, 80)
        
        # 数字圆圈
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(box["x"] + 3.5), Inches(box["y"] + 0.2), Inches(0.8), Inches(0.8)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = box["color"]
        circle.line.fill.background()
        
        # 数字
        num_text = slide.shapes.add_textbox(
            Inches(box["x"] + 3.65), Inches(box["y"] + 0.35), Inches(0.5), Inches(0.5)
        )
        num_frame = num_text.text_frame
        num_para = num_frame.paragraphs[0]
        num_para.text = str(box["num"])
        num_para.font.size = Pt(24)
        num_para.font.bold = True
        num_para.font.color.rgb = RGBColor(255, 255, 255)
        num_para.alignment = PP_ALIGN.CENTER
    
    # 底部红色装饰线
    bottom_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5), Inches(7.2), Inches(12.5), Inches(0.2)
    )
    bottom_line.fill.solid()
    bottom_line.fill.fore_color.rgb = RGBColor(231, 76, 60)
    bottom_line.line.fill.background()
    
    # 保存文件
    output_path = "/home/admin/.openclaw/workspace/三课展示促能.pptx"
    prs.save(output_path)
    print(f"PPT 文件已创建：{output_path}")
    return output_path

if __name__ == "__main__":
    create_ppt()
