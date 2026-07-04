#!/usr/bin/env python3
"""
吉林医药学院 PPT 模板 - 绿色答辩版
参考风格：绿色主题、学术答辩、清新简约
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# 绿色主题配色方案
PRIMARY_GREEN   = RGBColor(0x2E, 0x7D, 0x32)   # 主绿色（深绿）
LIGHT_GREEN     = RGBColor(0x4C, 0xAF, 0x50)   # 亮绿色
PALE_GREEN      = RGBColor(0xE8, 0xF5, 0xE9)   # 浅绿背景
ACCENT_GOLD     = RGBColor(0xC9, 0xA9, 0x59)   # 金色点缀
DARK_TEXT       = RGBColor(0x33, 0x33, 0x33)   # 深色文字
WHITE           = RGBColor(0xFF, 0xFF, 0xFF)
GRAY            = RGBColor(0x88, 0x88, 0x88)
LIGHT_GRAY      = RGBColor(0xF5, 0xF5, 0xF5)

W = Inches

def add_shape(slide, shape_type, left, top, width, height, fill_color=None, line_color=None, line_width=1):
    s = slide.shapes.add_shape(shape_type, left, top, width, height)
    if fill_color:
        s.fill.solid()
        s.fill.fore_color.rgb = fill_color
    else:
        s.fill.background()
    if line_color:
        s.line.color.rgb = line_color
        s.line.width = Pt(line_width)
    else:
        s.line.fill.background()
    return s

def add_textbox(slide, left, top, width, height, text, font_size=18, bold=False, color=DARK_TEXT, alignment=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    return tb

def create_ppt():
    prs = Presentation()
    prs.slide_width = W(13.333)
    prs.slide_height = W(7.5)
    
    # ============ 第 1 页：封面（绿色主题）============
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 绿色背景
    bg = add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, W(13.333), W(7.5), PRIMARY_GREEN)
    
    # 浅绿色装饰区域（左上角）
    dec1 = add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, W(5), W(3), LIGHT_GREEN)
    dec1.fill.transparency = 0.3
    
    # 白色云朵装饰（底部）
    cloud = add_shape(slide, MSO_SHAPE.CLOUD, W(1), W(5.5), W(11.333), W(2), WHITE)
    cloud.fill.transparency = 0.8
    
    # 校徽占位圆
    logo_circle = add_shape(slide, MSO_SHAPE.OVAL, W(10.5), W(0.5), W(1.5), W(1.5), WHITE)
    logo_text = add_textbox(slide, W(10.7), W(0.8), W(1.1), W(0.8), '校徽', 14, True, PRIMARY_GREEN, PP_ALIGN.CENTER)
    
    # 英文标题
    add_textbox(slide, W(1.5), W(1.2), W(8), W(0.6), 'ACADEMIC TEMPLATE', 16, False, WHITE, PP_ALIGN.LEFT)
    
    # 主标题
    add_textbox(slide, W(1.5), W(2.0), W(9), W(1.2), '吉林医药学院', 44, True, WHITE, PP_ALIGN.LEFT)
    add_textbox(slide, W(1.5), W(3.2), W(9), W(1.2), '毕业答辩', 44, True, WHITE, PP_ALIGN.LEFT)
    
    # 分隔线
    add_shape(slide, MSO_SHAPE.RECTANGLE, W(1.5), W(4.3), W(2), W(0.04), WHITE)
    
    # 汇报人信息
    info_box = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, W(1.5), W(4.6), W(6), W(0.8), WHITE)
    tf = info_box.text_frame
    p = tf.paragraphs[0]
    p.text = '汇报人：XXX    指导教师：XXX'
    p.font.size = Pt(16)
    p.font.color.rgb = PRIMARY_GREEN
    p.alignment = PP_ALIGN.CENTER
    
    # 手持论文图片占位区（右侧）
    photo_frame = add_shape(slide, MSO_SHAPE.RECTANGLE, W(8.5), W(2.5), W(3.5), W(4.5), WHITE)
    photo_frame.line.color.rgb = WHITE
    photo_frame.line.width = Pt(2)
    add_textbox(slide, W(8.8), W(4.2), W(2.9), W(1.0), '手持论文\n图片占位', 14, False, GRAY, PP_ALIGN.CENTER)
    
    # ============ 第 2 页：目录 ============
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 白色背景
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, W(13.333), W(7.5), WHITE)
    
    # 顶部绿色条
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, W(13.333), W(0.15), PRIMARY_GREEN)
    
    # 页面标题
    add_textbox(slide, W(1), W(0.6), W(10), W(0.8), '目录', 40, True, PRIMARY_GREEN, PP_ALIGN.LEFT)
    
    # 目录项
    toc_items = [
        ('01', '研究背景与意义'),
        ('02', '国内外研究现状'),
        ('03', '研究方法与内容'),
        ('04', '研究结果与分析'),
        ('05', '结论与展望'),
    ]
    
    for i, (num, text) in enumerate(toc_items):
        y = W(2.0 + i * 0.9)
        
        # 序号框
        num_box = add_shape(slide, MSO_SHAPE.RECTANGLE, W(1.2), y, W(1.2), W(0.7), PRIMARY_GREEN)
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = num
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        
        # 目录文字
        add_textbox(slide, W(2.8), y, W(8), W(0.7), text, 24, False, DARK_TEXT)
    
    # 页脚装饰
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0, W(7.35), W(13.333), W(0.15), LIGHT_GREEN)
    
    # ============ 第 3 页：内容页（图文）============
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 白色背景
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, W(13.333), W(7.5), WHITE)
    
    # 顶部绿色条
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, W(13.333), W(0.15), PRIMARY_GREEN)
    
    # 章节标题
    add_textbox(slide, W(1), W(0.5), W(10), W(0.8), '2.1 国内外研究现状', 32, True, PRIMARY_GREEN, PP_ALIGN.LEFT)
    
    # 左侧文本框
    left_text = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, W(0.8), W(1.6), W(5.5), W(5.2), PALE_GREEN)
    tf = left_text.text_frame
    tf.word_wrap = True
    
    content = [
        ('这里是小标题', True, 22, PRIMARY_GREEN),
        ('', False, 8, DARK_TEXT),
        ('这是正文文字内容这是正文文字内容这是正文文字内容这是正文文字内容这是正文文字内容这是正文文字内容这是正文文字内容这是正文文字内容', False, 14, DARK_TEXT),
    ]
    
    for i, (text, bold, size, color) in enumerate(content):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.space_after = Pt(6)
    
    # 中间图表占位
    chart_box = add_shape(slide, MSO_SHAPE.RECTANGLE, W(6.6), W(1.6), W(3.5), W(3.5), PALE_GREEN)
    chart_box.line.color.rgb = PRIMARY_GREEN
    chart_box.line.width = Pt(1.5)
    add_textbox(slide, W(6.8), W(2.8), W(3.1), W(1.5), '📊\n柱状图占位', 18, False, PRIMARY_GREEN, PP_ALIGN.CENTER)
    
    # 右侧文本框
    right_text = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, W(6.6), W(5.3), W(3.5), W(1.5), PALE_GREEN)
    tf = right_text.text_frame
    p = tf.paragraphs[0]
    p.text = '相关图表描述\n文字说明内容'
    p.font.size = Pt(14)
    p.font.color.rgb = DARK_TEXT
    
    # 右下角图片占位
    img_box = add_shape(slide, MSO_SHAPE.RECTANGLE, W(10.3), W(1.6), W(2.2), W(5.2), LIGHT_GRAY)
    add_textbox(slide, W(10.5), W(3.8), W(1.8), W(1.0), '图片\n占位', 14, False, GRAY, PP_ALIGN.CENTER)
    
    # ============ 第 4 页：章节标题页 ============
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 绿色背景
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, W(13.333), W(7.5), PRIMARY_GREEN)
    
    # 左侧装饰
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, W(0.3), W(7.5), LIGHT_GREEN)
    
    # 章节序号
    add_textbox(slide, W(1.5), W(2.0), W(2), W(1.5), '第一部分', 28, False, WHITE, PP_ALIGN.LEFT)
    
    # 章节标题
    add_textbox(slide, W(1.5), W(3.2), W(10), W(1.5), '这里是章节标题 1', 48, True, WHITE, PP_ALIGN.LEFT)
    
    # 装饰线
    add_shape(slide, MSO_SHAPE.RECTANGLE, W(1.5), W(4.8), W(3), W(0.04), LIGHT_GREEN)
    
    # 右侧手持教鞭图片占位
    pointer_img = add_shape(slide, MSO_SHAPE.RECTANGLE, W(9), W(2.5), W(3.5), W(4), WHITE)
    pointer_img.fill.transparency = 0.7
    add_textbox(slide, W(9.5), W(4.0), W(2.5), W(1.0), '手持教鞭\n图片占位', 14, False, GRAY, PP_ALIGN.CENTER)
    
    # ============ 第 5 页：结束页 ============
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 绿色背景
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, W(13.333), W(7.5), PRIMARY_GREEN)
    
    # 白色云朵装饰（底部）
    cloud = add_shape(slide, MSO_SHAPE.CLOUD, W(2), W(5.5), W(9.333), W(2), WHITE)
    cloud.fill.transparency = 0.8
    
    # 感谢文字
    add_textbox(slide, W(3), W(2.5), W(7.333), W(1.5), '感谢聆听', 48, True, WHITE, PP_ALIGN.CENTER)
    add_textbox(slide, W(3.5), W(4.0), W(6.333), W(0.8), '敬请各位专家批评指正', 24, False, WHITE, PP_ALIGN.CENTER)
    
    # 学校名称
    add_textbox(slide, W(4), W(6.5), W(5.333), W(0.6), '吉林医药学院', 20, False, RGBColor(0xCC, 0xDD, 0xCC), PP_ALIGN.CENTER)
    
    # 保存
    output = '/home/admin/.openclaw/workspace/吉林医药学院_PPT 模板_绿色答辩版.pptx'
    prs.save(output)
    print(f'✅ 绿色答辩版 PPT 模板已生成：{output}')
    return output

if __name__ == '__main__':
    create_ppt()
