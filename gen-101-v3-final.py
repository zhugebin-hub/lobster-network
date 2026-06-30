#!/usr/bin/env python3
"""
"101计划"首批核心课程培育推进会汇报PPT
融合：第一次研讨会纪要 + AI原生教学范式教改论文
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

C = {
    'primary': RGBColor(0x1A, 0x36, 0x5D),
    'secondary': RGBColor(0x2B, 0x6C, 0xB0),
    'accent': RGBColor(0xE8, 0x83, 0x2A),
    'light': RGBColor(0xF7, 0xF9, 0xFC),
    'white': RGBColor(0xFF, 0xFF, 0xFF),
    'text': RGBColor(0x2D, 0x37, 0x48),
    'muted': RGBColor(0x71, 0x80, 0x96),
    'success': RGBColor(0x38, 0xA1, 0x69),
    'teal': RGBColor(0x31, 0x97, 0x95),
}

FONT_NAME = 'Microsoft YaHei'

def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, x, y, w, h, fill_color, line_color=None, lw=0, radius=False):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(lw)
    else:
        shape.line.fill.background()
    return shape

def tb(slide, x, y, w, h, text, fs=14, color=C['text'], bold=False, align=PP_ALIGN.LEFT, italic=False):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(fs)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.italic = italic
    p.font.name = FONT_NAME
    p.alignment = align
    return txBox

def mtext(slide, x, y, w, h, lines, fs=14, color=C['text'], bold=False, align=PP_ALIGN.LEFT, italic=False, sp=2):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(fs)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.italic = italic
        p.font.name = FONT_NAME
        p.alignment = align
        p.space_after = Pt(sp)
    return txBox

def header(slide, title):
    add_rect(slide, 0, 0, 13.33, 1.1, C['primary'])
    tb(slide, 0.5, 0.15, 12, 0.9, title, fs=26, color=C['white'], bold=True)

# ============================================================
# SLIDE 1: Cover
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, C['primary'])
add_rect(s, 0, 0, 13.33, 0.12, C['accent'])
tb(s, 1, 1.0, 11.33, 1.0,
   '\u6559\u80b2\u90e8\u201c101\u8ba1\u5212\u201d\u9996\u6279\u6838\u5fc3\u8bfe\u7a0b\u57f9\u80b2\u63a8\u8fdb\u4f1a',
   fs=34, color=C['white'], bold=True, align=PP_ALIGN.CENTER)
tb(s, 1, 2.3, 11.33, 0.8,
   '\u57fa\u4e8e AI \u539f\u751f\u6559\u5b66\u8303\u5f0f\u7684\u8ba1\u7b97\u673a\u7f51\u7edc\u8bfe\u7a0b\u6539\u9769\u5b9e\u8df5',
   fs=24, color=C['accent'], bold=True, align=PP_ALIGN.CENTER)
mtext(s, 1, 5.5, 11.33, 0.8,
      ['\u6d59\u6c5f\u5de5\u5546\u5927\u5b66 \u00b7 \u8bf8\u845b\u658c\u56e2\u961f',
       '2026\u5e746\u670825\u65e5 | \u79d1\u521b\u5927\u697c206'],
      fs=15, color=C['muted'], align=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 2: 101 Plan Background
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, C['white'])
header(s, '\u201c101\u8ba1\u5212\u201d\u80cc\u666f\u4e0e\u9879\u76ee\u5b9a\u4f4d')

tb(s, 0.8, 1.4, 3, 0.5, '\u653f\u7b56\u80cc\u666f', fs=20, color=C['primary'], bold=True)
mtext(s, 0.8, 2.0, 5.5, 2.5,
      ['\u2022 \u6559\u80b2\u90e8\u7edf\u7b79\u7684\u62d4\u5c16\u521b\u65b0\u4eba\u624d\u57f9\u517b\u7b51\u57fa\u6027\u5de5\u7a0b',
       '\u2022 \u6c47\u805a\u9876\u5c16\u9ad8\u6821\u3001\u9876\u5c16\u5e08\u8d44\u3001\u9876\u5c16\u51fa\u7248\u5355\u4f4d',
       '\u2022 \u4ee5\u8bfe\u7a0b\u3001\u6559\u6750\u3001\u6559\u5e08\u548c\u5b9e\u8df5\u9879\u76ee\u4e3a\u6838\u5fc3\u8981\u7d20',
       '\u2022 \u8ba1\u7b97\u673a\u79d1\u5b66\u4e0e\u6280\u672f\u9886\u57df\u9996\u6279\u6838\u5fc3\u8bfe\u7a0b'],
      fs=15, color=C['text'])

add_rect(s, 7, 1.4, 5.5, 2.5, C['light'], C['accent'], 2, True)
tb(s, 7.2, 1.5, 5, 0.5, '\u9879\u76ee\u5b9a\u4f4d', fs=20, color=C['primary'], bold=True, align=PP_ALIGN.CENTER)
mtext(s, 7.3, 2.0, 5, 1.8,
      ['\u6a21\u5757\u4e00\uff1a\u8ba1\u7b97\u673a\u7f51\u7edc\uff08\u672c\u79d1\u00b7\u4e13\u4e1a\u6838\u5fc3\u8bfe\uff09',
       '\u5efa\u8bbe\u5468\u671f\uff1a2026.01 - 2027.12\uff082\u5e74\uff09',
       '\u6838\u5fc3\u6539\u9769\uff1a\u56fd\u4ea7\u4e91\u5e73\u53f0\u4e3a\u5e95\u5ea7\uff0c\u6df1\u5ea6\u878d\u5408\u667a\u80fd\u4f53\u5de5\u5177\uff0c',
       '\u6784\u5efa\u201c\u667a\u80fd\u751f\u6210-\u79c1\u6709\u7b54\u7591-\u667a\u6167\u7ba1\u7406-\u5b9e\u6218\u90e8\u7f72\u201d\u5168\u94fe\u6761\u6559\u5b66\u65b0\u8303\u5f0f'],
      fs=14, color=C['text'])

add_rect(s, 0.8, 4.3, 11.5, 2.8, C['primary'], radius=True)
tb(s, 1, 4.4, 3, 0.5, '\u5b98\u65b9\u8fdb\u5ea6\u5b89\u6392', fs=18, color=C['accent'], bold=True)
mtext(s, 1, 4.95, 11, 2.0,
      ['\u2022 2026.1-12\uff1a\u8bfe\u7a0b\u5efa\u8bbe\u9636\u6bb5\uff08\u77e5\u8bc6\u4f53\u7cfb\u68b3\u7406\u3001\u6559\u6750\u64b0\u5199\u3001\u5b9e\u8df5\u8d44\u6e90\u5efa\u8bbe\uff09',
       '\u2022 2027.1-12\uff1a\u8bd5\u70b9\u4e0e\u5b8c\u5584\u9636\u6bb5\uff08\u4e0d\u5c11\u4e8e5\u6240\u9ad8\u6821\u8bd5\u7528\uff0c\u6839\u636e\u53cd\u9988\u8fed\u4ee3\uff09',
       '\u2022 2026.11-12\uff1a\u4e2d\u671f\u5de5\u4f5c\u603b\u7ed3\uff0c\u786e\u7acb\u201c\u56db\u4e2a\u6838\u5fc3\u201d\u5efa\u8bbe\u8d28\u91cf\u6807\u51c6\u4e0e\u8ba4\u5b9a\u529e\u6cd5'],
      fs=14, color=C['white'])

# ============================================================
# SLIDE 3: Team Foundation
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, C['white'])
header(s, '\u56e2\u961f\u73b0\u6709\u6210\u679c\u57fa\u7840')

ach = [
    { 't': 'MOOC\u5728\u7ebf\u8bfe\u7a0b',
      'd': '\u300a\u9ad8\u7ea7\u7f51\u7edc\u901a\u4fe1\u539f\u7406\u5b9e\u8df5\u300b\n\u4e2d\u56fd\u5927\u5b66MOOC | 228\u4eba\u5b66\u4e60\n18\u5468\u8bfe\u7a0b\uff08\u5df2\u5b8c\u621014\u5468\uff09\n\u5408\u4f5c\u4f01\u4e1a\uff1a\u963f\u91cc\u4e91', 'c': C['secondary'] },
    { 't': '\u5df2\u51fa\u7248\u6559\u6750',
      'd': '\u300a\u7f51\u7edc\u901a\u4fe1\u539f\u7406\u5b9e\u8df5\u300b\n\uff08\u6e05\u534e\u51fa\u7248\u793e,2024\uff09\n\u300a\u7cfb\u7edf\u7ea7\u7f16\u7a0b\u53ca\u5206\u5e03\u5f0f\u5e94\u7528\u5b9e\u73b0\u6280\u672f\u300b\n\uff08\u6e05\u534e\u51fa\u7248\u793e,\u5df2\u51fa\u7248\uff09', 'c': C['teal'] },
    { 't': '\u7701\u7ea7\u4e00\u6d41\u5b9e\u9a8c\u8bfe',
      'd': '\u300a\u8ba1\u7b97\u673a\u7f51\u7edc\u5b9e\u9a8c\u300b\n7\u671f\u8fd0\u884c | 623\u4eba\u6b21\u9009\u8bfe\n22\u6240\u9ad8\u6821\u8986\u76d6\n\u7d2f\u8ba1\u8bbf\u95ee59.4\u4e07\u6b21', 'c': C['success'] },
    { 't': '\u5b9e\u9a8c\u5e73\u53f0',
      'd': '\u963f\u91cc\u4e91\u4e91\u5b9e\u9a8c\u5ba4\u5e73\u53f0\nMininet / OpenDaylight / OpenStack\n\u652f\u6301\u7f51\u7edc\u865a\u62df\u5316\u5b9e\u6218\n\u4f01\u4e1a\u7ea7\u5b9e\u9a8c\u73af\u5883', 'c': C['accent'] }
]

ax = 0.5
for a in ach:
    add_rect(s, ax, 1.5, 3.05, 4.0, C['light'], a['c'], 2, True)
    tb(s, ax + 0.15, 1.6, 2.75, 0.5, a['t'], fs=17, color=a['c'], bold=True, align=PP_ALIGN.CENTER)
    mtext(s, ax + 0.15, 2.2, 2.75, 3.0, a['d'].split('\n'), fs=13, color=C['text'])
    ax += 3.2

tb(s, 0.8, 5.7, 11, 0.5,
   '\u56e2\u961f\u6838\u5fc3\u6210\u5458\uff1a\u8bf8\u845b\u658c\uff08\u6559\u6388\uff09\u3001\u91d1\u84c9\uff08\u526f\u6559\u6388\uff09\u3001\u9ad8\u660e\u3001\u674e\u4f20\u714c\u3001\u848b\u732e',
   fs=14, color=C['muted'], align=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 4: Core Reform Direction (AI Native Paradigm)
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, C['white'])
header(s, '\u6838\u5fc3\u6539\u9769\u65b9\u5411\uff1aAI\u539f\u751f\u6559\u5b66\u65b0\u8303\u5f0f')

tb(s, 0.8, 1.4, 3, 0.5, '\u7406\u5ff5\u8f6c\u53d8', fs=20, color=C['accent'], bold=True)
mtext(s, 0.8, 2.0, 5.5, 2.8,
      ['\u2022 \u4ece\u201c\u77e5\u8bc6\u704c\u8f93\u201d \u2192 \u201cAI\u5e94\u7528\u7684\u5b9e\u8df5\u573a\u666f\u201d',
       '\u2022 \u201c\u8bfe\u7a0b\u4e0d\u518d\u662f\u6559\u5b66\u6838\u5fc3\uff0c\u800c\u662fAI\u5de5\u5177\u5e94\u7528\u7684\u80cc\u666f\u573a\u666f\u201d',
       '\u2022 \u6838\u5fc3\u76ee\u6807\uff1a\u57f9\u517b\u5bf9AI\u751f\u6210\u5185\u5bb9\u7684\u5224\u65ad\u529b\u4e0e\u6279\u5224\u6027\u601d\u7ef4',
       '\u2022 \u4e94\u5e74\u540e\u4e13\u4e1a\u6559\u80b2\u5c06\u5168\u9762\u91cd\u6784',
       '\u2022 \u6559\u5e08\u9700\u5148\u76f8\u4fe1AI\u7684\u6f5c\u529b\uff0c\u624d\u80fd\u8c03\u6574\u6559\u5b66\u7b56\u7565'],
      fs=14, color=C['text'])

add_rect(s, 7, 1.3, 5.5, 3.5, C['primary'], radius=True)
tb(s, 7.2, 1.4, 5, 0.5, '\u5168\u94fe\u6761\u6559\u5b66\u65b0\u8303\u5f0f', fs=20, color=C['accent'], bold=True, align=PP_ALIGN.CENTER)
mtext(s, 7.3, 2.1, 5, 2.5,
      ['\u2460 \u667a\u80fd\u751f\u6210 \u2014 AI\u8f85\u52a9\u77e5\u8bc6\u70b9\u5185\u5bb9\u751f\u6210',
       '\u2461 \u79c1\u6709\u7b54\u7591 \u2014 \u77e5\u8bc6\u5e93\u673a\u5668\u4eba24\u5c0f\u65f6\u7cbe\u51c6\u7b54\u7591',
       '\u2462 \u667a\u6167\u7ba1\u7406 \u2014 MOOC\u6155\u8bfe\u5802+\u9489\u9489AI\u52a9\u7406\u5168\u6d41\u7a0b\u8986\u76d6',
       '\u2463 \u5b9e\u6218\u90e8\u7f72 \u2014 \u56fd\u4ea7\u4e91\u5e73\u53f0\u771f\u5b9e\u73af\u5883\u9a8c\u8bc1'],
      fs=15, color=C['white'])

add_rect(s, 0.8, 5.2, 11.5, 1.3, C['light'], C['accent'], 2, True)
mtext(s, 1, 5.25, 11, 1.2,
      ['\u201c\u672a\u6765\u8bfe\u7a0b\u5efa\u8bbe\u7684\u5173\u952e\u4e0d\u518d\u662f\u77e5\u8bc6\u4f20\u6388\uff0c\u800c\u662f\u57f9\u517b\u5b66\u751f\u5bf9AI\u751f\u6210\u5185\u5bb9\u7684\u5224\u65ad\u529b\u3002\u201d',
       '\u201c\u4ea4\u4e00\u767e\u4efd\u4f5c\u4e1a\uff0c\u53ea\u8981\u6709\u4e24\u4efd\u8d85\u51fa\u60f3\u8c61\uff0c\u5c31\u662f\u6210\u529f\u3002\u201d'],
      fs=15, color=C['primary'], italic=True, align=PP_ALIGN.CENTER)
tb(s, 1, 6.2, 11, 0.3, '\u2014\u2014 \u8bf8\u845b\u658c', fs=13, color=C['muted'], align=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 5: Five Reform Initiatives
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, C['white'])
header(s, '\u4e94\u5927\u6539\u9769\u4e3e\u63aa')

initiatives = [
    { 'n': '01', 't': '\u534f\u8bae\u673a\u5236\u53ef\u89c6\u5316',
      'd': '\u5229\u7528AI\u5de5\u5177\u751f\u6210\u52a8\u6001\u52a8\u753b\uff0c\u8986\u76d6STP\u3001\u8def\u7531\u67e5\u627e\u3001CHAP\u3001ACL\u3001NAT\u3001TCP\u4e09\u6b21\u63e1\u624b\u3001OpenFlow\u7b49\u6838\u5fc3\u534f\u8bae\u3002\u5b66\u751f\u4e3b\u5bfc\u751f\u6210\uff0c\u6559\u5e08\u7b5b\u9009\u6700\u4f18\uff0c\u5f62\u6210\u534f\u8bae\u52a8\u753b\u5e93\u3002',
      'c': C['secondary'] },
    { 'n': '02', 't': '\u8bfe\u7a0b\u79c1\u6709\u77e5\u8bc6\u5e93\u673a\u5668\u4eba',
      'd': '\u6574\u5408\u6559\u6750\u3001\u8bfe\u4ef6\u3001\u4e60\u9898\u3001FAQ\uff0c\u63d0\u4f9b24\u5c0f\u65f6\u4e2a\u6027\u5316\u8f85\u5bfc\u7b54\u7591\u3002\u6280\u672f\u5b9e\u73b0\uff1a\u767e\u70bc\u77e5\u8bc6\u5e93\u00d7OpenClaw\u96c6\u6210\uff0c\u4e09\u7ea7\u8bb0\u5fc6\u67b6\u6784\uff08L1\u5de5\u4f5c/L2\u77e5\u8bc6/L3\u957f\u671f\uff09\u3002',
      'c': C['teal'] },
    { 'n': '03', 't': '\u5168\u6d41\u7a0b\u667a\u6167\u6559\u5b66\u7ba1\u7406',
      'd': 'MOOC\u6155\u8bfe\u5802+\u9489\u9489AI\u52a9\u7406\uff0c\u8bfe\u524d-\u8bfe\u4e2d-\u8bfe\u540e\u5168\u6d41\u7a0b\u8986\u76d6\uff0c\u6570\u636e\u9a71\u52a8\u7684\u6559\u5b66\u6539\u8fdb\u3002\u901a\u8fc7\u9489\u9489\u7fa4\u4e0eAI\u52a9\u624b\u534f\u4f5c\uff0c\u5b9e\u73b0\u4f1a\u8bae\u7eaa\u8981\u81ea\u52a8\u751f\u6210\u3001\u8bfe\u7a0b\u603b\u7ed3\u667a\u80fd\u6c47\u603b\u3002',
      'c': C['accent'] },
    { 'n': '04', 't': '\u667a\u80fd\u4f53\u8f85\u52a9\u5b9e\u9a8c\u65b0\u8303\u5f0f',
      'd': 'AI\u8f85\u52a9\u751f\u6210\u547d\u4ee4+\u4eba\u5de5\u4e13\u4e1a\u6392\u9519+\u771f\u5b9e\u8bbe\u5907\u90e8\u7f72\u3002\u4fdd\u7559\u771f\u5b9e\u4ea4\u4e92\u903b\u8f91\uff0cAI\u8f85\u52a9\u6982\u5ff5\u7406\u89e3\u3002\u9a8c\u8bc1\u6027/\u8bbe\u8ba1\u6027/\u7efc\u5408\u6027=3:4:3\uff0c\u5b9e\u9a8c\u6548\u7387\u63d0\u534740%\u3002',
      'c': C['success'] },
    { 'n': '05', 't': '\u667a\u80fd\u4f53\u9a71\u52a8\u7684\u8bfe\u7a0b\u95e8\u6237\u5efa\u8bbe',
      'd': '\u8bfe\u7a0b\u5b98\u7f51\uff0c\u96c6\u6210\u6559\u5b66\u8d44\u6e90+\u5b9e\u9a8c\u6848\u4f8b+\u4f18\u79c0\u4f5c\u54c1+\u667a\u80fd\u4e92\u52a8\u3002\u5a92\u4f53\u5316\u8868\u8fbe\uff0c\u8d85\u8d8a\u4f20\u7edfPPT\u548c\u6559\u6750\u7684\u89c6\u89c9\u5448\u73b0\u9ad8\u5ea6\u3002',
      'c': RGBColor(0x9F, 0x7A, 0xEA) }
]

iy = 1.3
for init in initiatives:
    add_rect(s, 0.6, iy, 0.7, 1.0, init['c'], radius=True)
    tb(s, 0.65, iy + 0.15, 0.6, 0.7, init['n'], fs=20, color=C['white'], bold=True, align=PP_ALIGN.CENTER)
    tb(s, 1.5, iy + 0.05, 4.5, 0.4, init['t'], fs=16, color=init['c'], bold=True)
    tb(s, 1.5, iy + 0.45, 11.0, 0.55, init['d'], fs=12, color=C['text'])
    iy += 1.1

# ============================================================
# SLIDE 6: Student Assignment = Course Resource
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, C['white'])
header(s, '\u521b\u65b0\u6a21\u5f0f\uff1a\u201c\u5b66\u751f\u4f5c\u4e1a\u5373\u8bfe\u7a0b\u8d44\u6e90\u201d')

# Left: Process
tb(s, 0.8, 1.4, 4, 0.5, '\u64cd\u4f5c\u6d41\u7a0b', fs=20, color=C['accent'], bold=True)

steps = [
    ('\u2460 \u77e5\u8bc6\u70b9\u62c6\u89e3', '\u5c06\u8bfe\u7a0b\u77e5\u8bc6\u70b9\u62c6\u89e3\u4e3a\u5b66\u751f\u4f5c\u4e1a\u4efb\u52a1\uff0c\u6bcf\u8282\u8bfe1-2\u4e2a\u77e5\u8bc6\u70b9'),
    ('\u2461 \u5b66\u751fAI\u5b9e\u8df5', '\u5b66\u751f\u4f7f\u7528AI\u5de5\u5177\u751f\u6210\u53ef\u89c6\u5316\u5185\u5bb9\uff08\u52a8\u753b\u3001\u8bb2\u89e3\u3001\u4ee3\u7801\u7b49\uff09'),
    ('\u2462 \u6548\u679c\u7b5b\u9009', '\u6559\u5e08\u4ece\u5168\u73ed\u63d0\u4ea4\u7684AI\u751f\u6210\u5185\u5bb9\u4e2d\u6311\u9009\u6bcf\u77e5\u8bc6\u70b93\u4e2a\u6700\u4f18\u4f5c\u54c1'),
    ('\u2463 \u8d44\u6e90\u5e93\u5efa\u8bbe', '\u5f62\u6210\u53ef\u590d\u7528\u7684\u77e5\u8bc6\u70b9\u8d44\u6e90\u5e93\uff0c\u6bcf\u77e5\u8bc6\u70b9\u914d\u4e00\u4e2a\u4f18\u8d28AI\u751f\u6210\u5730\u5740')
]

sy = 2.0
for title, desc in steps:
    add_rect(s, 0.8, sy, 5.5, 0.9, C['light'], radius=True)
    tb(s, 1.0, sy + 0.05, 5.0, 0.35, title, fs=14, color=C['primary'], bold=True)
    tb(s, 1.0, sy + 0.4, 5.0, 0.45, desc, fs=11, color=C['text'])
    sy += 1.0

# Arrow between steps
for i in range(3):
    tb(s, 3.2, 2.55 + i*1.0, 0.8, 0.3, '\u2193', fs=16, color=C['accent'], align=PP_ALIGN.CENTER)

# Right: Evaluation shift
add_rect(s, 7, 1.3, 5.5, 3.8, C['light'], C['accent'], 2, True)
tb(s, 7.2, 1.4, 5, 0.5, '\u4ef7\u503c\u8bc4\u5224\u6743\u8f6c\u79fb', fs=18, color=C['primary'], bold=True, align=PP_ALIGN.CENTER)
mtext(s, 7.2, 2.0, 5.1, 3.0,
      ['\u5747\u503c\u5bfc\u5411 \u2192 \u5cf0\u503c\u53d1\u73b0',
      '  \u4e0d\u8981\u6c42\u6240\u6709\u4eba\u505a\u5bf9\uff0c\u800c\u662f\u53d1\u73b0\u4f18\u79c0\u4f5c\u54c1',
      '',
      '\u6559\u5e08\u4e3b\u89c2\u8bc4\u4ef7 \u2192 \u5b66\u751f\u4f5c\u54c1\u8d28\u91cf\u9a71\u52a8',
      '  \u4ee5\u4f5c\u54c1\u8d28\u91cf\u800c\u975e\u6559\u5e08\u504f\u597d\u4e3a\u6807\u51c6',
      '',
      '\u4e00\u6b21\u6027\u8bc4\u4ef7 \u2192 \u8fed\u4ee3\u4f18\u5316',
      '  \u4e0d\u8981\u6c42\u4e00\u6b21\u6027\u5b8c\u7f8e\uff0c\u901a\u8fc7\u6301\u7eed\u4f18\u5316\u63d0\u5347\u6548\u679c',
      '',
      '\u201c\u4e0d\u5728\u4e8eAI\u80fd\u5426\u5b8c\u6210\u4efb\u52a1\uff0c',
      '  \u800c\u5728\u4e8e\u82b1\u66f4\u591a\u65f6\u95f4\u540e\u662f\u5426\u80fd\u505a\u5f97\u66f4\u597d\u3002\u201d'],
      fs=13, color=C['text'])

# Bottom: SDP Case
add_rect(s, 0.8, 5.5, 11.5, 1.5, C['light'], C['secondary'], 1.5, True)
tb(s, 1, 5.55, 3, 0.4, 'SDP\u8bfe\u7a0b\u5b9e\u8df5\u6848\u4f8b', fs=16, color=C['secondary'], bold=True)
mtext(s, 1, 5.95, 11, 1.0,
      ['\u2022 \u77e5\u8bc6\u70b9\u62c6\u89e3\uff1a\u5c06SDP\u6838\u5fc3\u6982\u5ff5\uff08\u63a7\u5236\u5668\u3001\u6570\u636e\u5e73\u9762\u3001OpenFlow\u534f\u8bae\u7b49\uff09\u62c6\u89e3\u4e3a10\u4e2a\u4f5c\u4e1a\u4efb\u52a1',
       '\u2022 \u5b66\u751f\u5b9e\u8df5\uff1a\u4f7f\u7528AI\u5de5\u5177\u751f\u6210\u52a8\u753b\u3001\u8bb2\u89e3\u89c6\u9891\u3001\u4ee3\u7801\u5b9e\u73b0',
       '\u2022 \u6548\u679c\u7b5b\u9009\uff1a\u4ece100\u4efd\u4f5c\u4e1a\u4e2d\u7b5b\u9009\u51fa30\u4efd\u4f18\u8d28\u4f5c\u54c1\uff0c\u5f62\u6210SDP\u77e5\u8bc6\u70b9\u8d44\u6e90\u5e93'],
      fs=13, color=C['text'])

# ============================================================
# SLIDE 7: Three Core Tasks
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, C['white'])
header(s, '\u4e09\u5927\u6838\u5fc3\u4efb\u52a1\u4f53\u7cfb')

tasks = [
    { 'num': '\u4efb\u52a1\u4e00', 'title': '\u8bfe\u7a0b\u5efa\u8bbe', 'lead': '\u8bf8\u845b\u658c',
      'items': ['\u77e5\u8bc6\u4f53\u7cfb\u68b3\u7406\uff0850-60\u4e2a\u77e5\u8bc6\u70b9\uff09',
                '\u77e5\u8bc6\u70b9\u56fe\u8c31\u53ef\u89c6\u5316\u6784\u5efa',
                '\u8bfe\u7a0b\u5efa\u8bbe\u6307\u5357\u7f16\u5199',
                '\u672c\u5730\u5316\u6559\u5b66\u5927\u7eb2\u8bbe\u8ba1\uff08\u56fd\u4ea7\u4e91+\u667a\u80fd\u4f53\uff09',
                '\u6559\u5b66\u8d44\u6e90\u5305\u5f00\u53d1\uff08\u52a8\u753b/\u8bfe\u4ef6/\u4e60\u9898\uff09',
                '\u5b9e\u9a8c\u8bfe\u7a0b2.0\u5347\u7ea7',
                '\u8bfe\u7a0b\u95e8\u6237\u7f51\u7ad9\u5efa\u8bbe'],
      'c': C['secondary'] },
    { 'num': '\u4efb\u52a1\u4e8c', 'title': '\u6559\u6750\u7f16\u5199', 'lead': '\u8bf8\u845b\u658c+\u9ad8\u660e',
      'items': ['\u57fa\u4e8e\u73b0\u6709\u6559\u6750\u5347\u7ea7\uff08\u975e\u5168\u65b0\u7f16\u5199\uff09',
                '\u6570\u5b57\u6559\u6750\u914d\u5957\u89c6\u9891/\u52a8\u753b/\u4ea4\u4e92',
                '\u6bcf\u7ae0\u914d\u5957\u6848\u4f8b\u22653\u4e2a',
                '\u4f01\u4e1a\u771f\u5b9e\u6848\u4f8b\u5360\u6bd4\u226530%',
                '\u6e05\u534e\u5927\u5b66\u51fa\u7248\u793e\u5408\u4f5c\u7eed\u7b7e',
                '\u6837\u7ae0\u64b0\u5199 + \u6559\u6750\u521d\u7a3f',
                '\u6839\u636e\u8bd5\u7528\u53cd\u9988\u8fed\u4ee3\u4fee\u8ba2'],
      'c': C['teal'] },
    { 'num': '\u4efb\u52a1\u4e09', 'title': '\u5b9e\u8df5\u6848\u4f8b', 'lead': '\u8bf8\u845b\u658c+\u848b\u732e',
      'items': ['\u9a8c\u8bc1\u6027/\u8bbe\u8ba1\u6027/\u7efc\u5408\u6027=3:4:3',
                '\u667a\u80fd\u4f53\u8f85\u52a9\u5b9e\u9a8c\uff08AI\u751f\u6210\u547d\u4ee4+\u4eba\u5de5\u6392\u9519\uff09',
                '\u963f\u91cc\u4e91\u5e73\u53f0+\u591a\u4e91\u652f\u6301\u5347\u7ea7',
                '\u62d3\u5c55\u534e\u4e3a/\u534e\u4e09\u7b49\u4f01\u4e1a\u5408\u4f5c',
                '\u5b9e\u9a8c\u9879\u76ee\u6e05\u5355\uff08\u226510\u4e2a\uff09',
                '\u5b9e\u8df5\u6848\u4f8b\u7fa4\u5efa\u8bbe\uff08\u22655\u4e2a\uff09',
                '\u4ea7\u5b66\u7814\u5408\u4f5c\u8d44\u6e90\uff08\u4f01\u4e1a\u771f\u5b9e\u9879\u76ee\u22652\u4e2a\uff09'],
      'c': C['accent'] }
]

tx = 0.4
for t in tasks:
    add_rect(s, tx, 1.35, 4.1, 5.6, C['light'], t['c'], 2, True)
    tb(s, tx+0.15, 1.4, 3.8, 0.5, f"{t['num']}\uff1a{t['title']}", fs=18, color=t['c'], bold=True, align=PP_ALIGN.CENTER)
    tb(s, tx+0.15, 1.9, 3.8, 0.4, f"\u8d1f\u8d23\u4eba\uff1a{t['lead']}", fs=13, color=C['muted'], align=PP_ALIGN.CENTER)
    mtext(s, tx+0.15, 2.3, 3.8, 4.5, t['items'], fs=12, color=C['text'])
    tx += 4.25

# ============================================================
# SLIDE 8: Milestones
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, C['white'])
header(s, '\u5173\u952e\u91cc\u7a0b\u7891\u4e0e\u65f6\u95f4\u8282\u70b9')

milestones = [
    { 'd': '2026.06.30', 'i': '\u77e5\u8bc6\u4f53\u7cfb\u521d\u7a3f+\u6559\u5b66\u5927\u7eb2\u672c\u5730\u5316+\u7b2c1\u6b21\u7814\u8ba8\u8bb0\u5f55', 'ok': True },
    { 'd': '2026.07.15', 'i': '\u6559\u6750\u6846\u67b6\u8bbe\u8ba1+\u77e5\u8bc6\u70b9\u56fe\u8c31\u53ef\u89c6\u5316\u5b8c\u6210', 'ok': False },
    { 'd': '2026.07.31', 'i': '\u6559\u6750\u6837\u7ae02\u7ae0+\u5b9e\u9a8c\u6848\u4f8b5\u4e2a+\u5b9e\u9a8c\u9879\u76ee\u6e05\u5355', 'ok': False },
    { 'd': '2026.08.31', 'i': '\u534f\u8bae\u52a8\u753b5\u4e2a+\u77e5\u8bc6\u5e93\u673a\u5668\u4ebaMVP+\u4f01\u4e1a\u5408\u4f5c\u610f\u5411', 'ok': False },
    { 'd': '2026.09.30', 'i': '\u5168\u90e8\u4ea4\u4ed8\u7269\u521d\u7a3f+\u6848\u4f8b\u8d44\u6e90\u5e93+\u76ee\u6807\u9ad8\u6821\u8054\u7edc', 'ok': False },
    { 'd': '2026.10.31', 'i': '\u6559\u6750\u521d\u7a3f+\u5b9e\u9a8c\u8bfe\u7a0b2.0+\u4ea7\u5b66\u7814\u8d44\u6e90\u5e93+\u4e2d\u671f\u6750\u6599', 'ok': False },
    { 'd': '2026.11-12', 'i': '\u4e2d\u671f\u5de5\u4f5c\u603b\u7ed3+\u201c\u56db\u4e2a\u6838\u5fc3\u201d\u8d28\u91cf\u6807\u51c6\u5236\u5b9a', 'ok': False },
]

my = 1.4
for m in milestones:
    dc = C['success'] if m['ok'] else C['light']
    dl = C['success'] if m['ok'] else C['muted']
    dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1), Inches(my), Inches(0.35), Inches(0.35))
    dot.fill.solid(); dot.fill.fore_color.rgb = dc; dot.line.color.rgb = dl; dot.line.width = Pt(1.5)
    tb(s, 1.5, my-0.02, 2.2, 0.4, m['d'], fs=16, color=C['primary'], bold=True)
    tb(s, 3.8, my-0.02, 8, 0.4, m['i'], fs=14, color=C['text'])
    st = '\u2713 \u5df2\u5b8c\u6210' if m['ok'] else '\u25cb \u63a8\u8fdb\u4e2d'
    sc = C['success'] if m['ok'] else C['muted']
    tb(s, 12, my-0.02, 1.2, 0.4, st, fs=12, color=sc, bold=True)
    my += 0.8

# ============================================================
# SLIDE 9: Wisdom Tree + Online Course
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, C['white'])
header(s, '\u667a\u6167\u6811\u5728\u7ebf\u8bfe\u7a0b\u90e8\u7f72\u8fdb\u5c55')

tb(s, 0.8, 1.4, 3, 0.5, '\u5e73\u53f0\u4ef7\u503c', fs=18, color=C['primary'], bold=True)
mtext(s, 0.8, 2.0, 5.5, 2.0,
      ['\u2022 \u56fd\u5185\u9886\u5148\u8de8\u6821\u5171\u4eab\u5728\u7ebf\u6559\u80b2\u5e73\u53f0',
       '\u2022 \u652f\u6301\u5927\u89c4\u6a21\u5728\u7ebf\u5f00\u653e\u8bfe\u7a0b\u8fd0\u8425',
       '\u2022 \u5b8c\u5584\u5b66\u60c5\u5206\u6790\u4e0e\u8fc7\u7a0b\u8bc4\u4ef7\u4f53\u7cfb',
       '\u2022 \u5b66\u5206\u4e92\u8ba4\uff0c\u6269\u5927\u8bfe\u7a0b\u8f90\u5c04\u9762'],
      fs=14, color=C['text'])

tb(s, 7, 1.4, 4, 0.5, '\u5f53\u524d\u90e8\u7f72\u8fdb\u5ea6', fs=18, color=C['primary'], bold=True)

dl = [
    ('\u2713 \u8bfe\u7a0b\u6846\u67b6\u642d\u5efa', C['success']),
    ('\u2713 \u6559\u5b66\u89c6\u9891\u4e0a\u4f20', C['success']),
    ('\u2713 \u7ae0\u8282\u6d4b\u9a8c\u914d\u7f6e', C['success']),
    ('\u27f3 \u8ba8\u8bba\u533a\u4e0e\u4e92\u52a8\u6a21\u5757\u8bbe\u7f6e', C['accent']),
    ('\u27f3 \u4f5c\u4e1a\u6279\u6539\u89c4\u5219\u914d\u7f6e', C['accent']),
    ('\u25cb AI\u8f85\u52a9\u7b54\u7591\u6a21\u5757\u63a5\u5165', C['muted']),
    ('\u25cb \u671f\u672b\u8003\u8bd5\u4e0e\u8bc4\u4ef7\u914d\u7f6e', C['muted']),
]

dy = 2.0
for txt, col in dl:
    tb(s, 7, dy, 5.5, 0.45, txt, fs=14, color=col, bold='\u2713' in txt)
    dy += 0.45

add_rect(s, 0.8, 5.5, 11.5, 0.8, C['light'], radius=True)
tb(s, 1, 5.55, 4, 0.3, '\u8bfe\u7a0b\u5efa\u8bbe\u6574\u4f53\u8fdb\u5ea6\uff1a\u7ea665%', fs=14, color=C['primary'], bold=True)
add_rect(s, 5.5, 5.65, 5.5, 0.3, C['success'], radius=True)

# ============================================================
# SLIDE 10: Pilot Universities + Collaboration
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, C['white'])
header(s, '\u8bd5\u70b9\u9ad8\u6821\u4e0e\u534f\u4f5c\u673a\u5236')

tb(s, 0.8, 1.4, 4, 0.5, '\u8bd5\u70b9\u9ad8\u6821\u540d\u5355\uff085\u6240\uff09', fs=18, color=C['primary'], bold=True)
unis = ['\u6d59\u6c5f\u5de5\u4e1a\u5927\u5b66', '\u676d\u5dde\u7535\u5b50\u79d1\u6280\u5927\u5b66', '\u6d59\u6c5f\u7406\u5de5\u5927\u5b66', '\u5b81\u6ce2\u5927\u5b66', '\u6d59\u6c5f\u5e08\u8303\u5927\u5b66']
uy = 2.1
for u in unis:
    add_rect(s, 0.8, uy, 5.5, 0.55, C['light'], radius=True)
    tb(s, 1.1, uy+0.05, 4, 0.45, f"\u2b50 {u}", fs=16, color=C['primary'], bold=True)
    uy += 0.7

tb(s, 7.5, 1.4, 3, 0.5, '\u8bd5\u7528\u5185\u5bb9\u4e0e\u53cd\u9988', fs=18, color=C['primary'], bold=True)
mtext(s, 7.5, 2.0, 5, 2.5,
      ['\u2022 \u914d\u5957\u6559\u6750 + \u5b9e\u9a8c\u6307\u5bfc\u4e66',
       '\u2022 \u5728\u7ebf\u8bfe\u7a0b + \u5b9e\u8df5\u6848\u4f8b',
       '\u2022 \u4e91\u5e73\u53f0\u5b9e\u9a8c\u73af\u5883 + AI\u7b54\u7591\u6a21\u5757',
       '',
       '\u53cd\u9988\u673a\u5236\uff1a',
       '\u95ee\u5377\u8c03\u67e5 + \u6df1\u5ea6\u8bbf\u8c08 + \u6570\u636e\u5206\u6790',
       '\u76ee\u6807\u6837\u672c\uff1a\u4e0d\u5c11\u4e8e500\u4efd',
       '\u8d1f\u8d23\u4eba\uff1a\u674e\u4f20\u714c | \u65f6\u95f4\uff1a2026.09-2027.06'],
      fs=14, color=C['text'])

# Collaboration table
add_rect(s, 0.8, 5.5, 11.5, 1.6, C['light'], radius=True)
tb(s, 1, 5.55, 3, 0.4, '\u534f\u4f5c\u673a\u5236', fs=16, color=C['primary'], bold=True)
mtext(s, 1, 5.95, 11, 1.0,
      ['\u2022 \u7ebf\u4e0a\u78b0\u5934\u4f1a\uff1a\u53cc\u5468 | \u5468\u4e94 15:00-16:00\uff08\u8bf8\u845b\u658c\uff09',
       '\u2022 \u7ebf\u4e0b\u7814\u8ba8\u4f1a\uff1a\u6bcf\u5b63\u5ea6\uff08\u5168\u4f53\uff09 | \u6587\u6863\u5f52\u6863\uff1a\u6bcf\u6708\u5e95\uff08\u91d1\u84c9\uff09',
       '\u2022 \u4e2d\u671f\u603b\u7ed3\uff1a2026\u5e7411-12\u6708\uff08\u674e\u4f20\u714c\uff09'],
      fs=13, color=C['text'])

# ============================================================
# SLIDE 11: Expected Outcomes
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, C['white'])
header(s, '\u9884\u671f\u6210\u679c\u4e0e\u5c55\u671b')

outcomes = [
    '\u5efa\u6210\u7b26\u5408\u201c101\u8ba1\u5212\u201d\u6807\u51c6\u7684\u8ba1\u7b97\u673a\u7f51\u7edc\u6838\u5fc3\u8bfe\u7a0b\u4f53\u7cfb',
    '\u51fa\u7248\u914d\u5957\u6570\u667a\u5316\u6559\u6750\uff08\u6e05\u534e\u5927\u5b66\u51fa\u7248\u793e\uff09\u4e0e\u534f\u8bae\u52a8\u753b\u5e93',
    '\u5f62\u6210\u8bfe\u7a0b\u79c1\u6709\u77e5\u8bc6\u5e93\u673a\u5668\u4eba + \u5168\u6d41\u7a0b\u667a\u6167\u6559\u5b66\u7ba1\u7406\u4f53\u7cfb',
    '\u5b8c\u6210\u667a\u6167\u6811\u5e73\u53f0\u5728\u7ebf\u8bfe\u7a0b\u90e8\u7f72\u4e0e\u8de8\u6821\u8fd0\u8425',
    '\u5b9e\u9a8c\u6548\u7387\u63d0\u534740%+\uff0c\u6bcf\u5b66\u671f\u53d7\u76ca200+\u5b66\u751f',
    '\u4e0d\u5c11\u4e8e5\u6240\u9ad8\u6821\u8bd5\u70b9\uff0c\u8986\u76d6500+\u5b66\u751f\u53cd\u9988\u6570\u636e',
    '\u6253\u9020\u201c101\u8ba1\u5212\u201d\u672c\u5730\u5316\u201cAI+\u6559\u80b2\u201d\u53ef\u590d\u5236\u3001\u53ef\u63a8\u5e7f\u5178\u578b\u6848\u4f8b',
]

oy = 1.4
for o in outcomes:
    add_rect(s, 0.8, oy, 0.55, 0.55, C['accent'], radius=True)
    tb(s, 1.55, oy-0.03, 10.5, 0.6, o, fs=16, color=C['text'], bold=True)
    oy += 0.8

# ============================================================
# SLIDE 12: Thank You
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, C['primary'])
add_rect(s, 0, 0, 13.33, 0.12, C['accent'])
tb(s, 1, 2.0, 11.33, 1.0,
   '\u611f\u8c22\u8046\u542c\uff01',
   fs=44, color=C['white'], bold=True, align=PP_ALIGN.CENTER)
tb(s, 1, 3.0, 11.33, 0.7,
   '\u6b22\u8fce\u6279\u8bc4\u6307\u6b63',
   fs=24, color=C['accent'], align=PP_ALIGN.CENTER)
mtext(s, 1, 5.5, 11.33, 0.8,
      ['\u8bf8\u845b\u658c\u56e2\u961f \u00b7 \u6d59\u6c5f\u5de5\u5546\u5927\u5b66',
       '\u8ba1\u7b97\u673a\u7f51\u7edc\u201c101\u8ba1\u5212\u201d\u8bfe\u7a0b\u5efa\u8bbe\u9879\u76ee'],
      fs=16, color=C['muted'], align=PP_ALIGN.CENTER)

output_path = '/home/admin/.openclaw/workspace/101-network-course-v3-final.pptx'
prs.save(output_path)
print(f'PPT saved: {output_path}')
