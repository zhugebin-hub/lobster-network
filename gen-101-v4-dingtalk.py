#!/usr/bin/env python3
"""
"101计划"首批核心课程培育推进会汇报PPT
融合：第一次研讨会 + AI原生教学范式教改论文 + 钉钉群AI原生管理机制
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

C = {
    'primary': RGBColor(0x1A, 0x36, 0x5D),
    'secondary': RGBColor(0x2B, 0x6C, 0xB0),
    'accent': RGBColor(0xE8, 0x83, 0x2A),
    'light': RGBColor(0xF7, 0xF9, 0xFC),
    'white': RGBColor(0xFF, 0xFF, 0xFF),
    'text': RGBColor(0x2D, 0x37, 0x48),
    'muted': RGBColor(0x71, 0x80, 0x96),
    'success': RGBColor(0x38, 0xA1, 0x69),
    'teal': RGBColor(0x31, 0x97, 0x95),
    'purple': RGBColor(0x9F, 0x7A, 0xEA),
}
FN = 'Microsoft YaHei'

def bg(s, c):
    f = s.background.fill; f.solid(); f.fore_color.rgb = c

def rect(s, x, y, w, h, fc, lc=None, lw=0, r=False):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if r else MSO_SHAPE.RECTANGLE,
                             Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fc
    if lc: sh.line.color.rgb = lc; sh.line.width = Pt(lw)
    else: sh.line.fill.background()
    return sh

def tb(s, x, y, w, h, t, fs=14, c=C['text'], b=False, a=PP_ALIGN.LEFT, i=False):
    bx = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = bx.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = t; p.font.size = Pt(fs); p.font.color.rgb = c
    p.font.bold = b; p.font.italic = i; p.font.name = FN; p.alignment = a
    return bx

def mt(s, x, y, w, h, lines, fs=14, c=C["text"], b=False, a=PP_ALIGN.LEFT, italic=False, sp=2):
    bx = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = bx.text_frame; tf.word_wrap = True
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = line; p.font.size = Pt(fs); p.font.color.rgb = c
        p.font.bold = b; p.font.italic = italic; p.font.name = FN; p.alignment = a
        p.space_after = Pt(sp)
    return bx

def hdr(s, title):
    rect(s, 0, 0, 13.33, 1.1, C['primary'])
    tb(s, 0.5, 0.15, 12, 0.9, title, fs=26, c=C['white'], b=True)

def add_img(slide, img_path, x, y, w, h):
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(x), Inches(y), Inches(w), Inches(h))
        return True
    return False

# ============================================================
# SLIDE 1: Cover
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, C['primary'])
rect(s, 0, 0, 13.33, 0.12, C['accent'])
tb(s, 1, 1.0, 11.33, 1.0,
   '\u6559\u80b2\u90e8\u201c101\u8ba1\u5212\u201d\u9996\u6279\u6838\u5fc3\u8bfe\u7a0b\u57f9\u80b2\u63a8\u8fdb\u4f1a',
   fs=34, c=C['white'], b=True, a=PP_ALIGN.CENTER)
tb(s, 1, 2.3, 11.33, 0.8,
   '\u57fa\u4e8eAI\u539f\u751f\u6559\u5b66\u8303\u5f0f\u7684\u8ba1\u7b97\u673a\u7f51\u7edc\u8bfe\u7a0b\u6539\u9769\u5b9e\u8df5',
   fs=24, c=C['accent'], b=True, a=PP_ALIGN.CENTER)
mt(s, 1, 5.5, 11.33, 0.8,
   ['\u6d59\u6c5f\u5de5\u5546\u5927\u5b66 \u00b7 \u8bf8\u845b\u658c\u56e2\u961f',
    '2026\u5e746\u670825\u65e5 | \u79d1\u521b\u5927\u697c206'],
   fs=15, c=C['muted'], a=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 2: 101 Plan Background
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, C['white'])
hdr(s, '\u201c101\u8ba1\u5212\u201d\u80cc\u666f\u4e0e\u9879\u76ee\u5b9a\u4f4d')

tb(s, 0.8, 1.4, 3, 0.5, '\u653f\u7b56\u80cc\u666f', fs=20, c=C['primary'], b=True)
mt(s, 0.8, 2.0, 5.5, 2.5,
   ['\u2022 \u6559\u80b2\u90e8\u7edf\u7b79\u7684\u62d4\u5c16\u521b\u65b0\u4eba\u624d\u57f9\u517b\u7b51\u57fa\u6027\u5de5\u7a0b',
    '\u2022 \u6c47\u805a\u9876\u5c16\u9ad8\u6821\u3001\u9876\u5c16\u5e08\u8d44\u3001\u9876\u5c16\u51fa\u7248\u5355\u4f4d',
    '\u2022 \u4ee5\u8bfe\u7a0b\u3001\u6559\u6750\u3001\u6559\u5e08\u548c\u5b9e\u8df5\u9879\u76ee\u4e3a\u6838\u5fc3\u8981\u7d20',
    '\u2022 \u8ba1\u7b97\u673a\u79d1\u5b66\u4e0e\u6280\u672f\u9886\u57df\u9996\u6279\u6838\u5fc3\u8bfe\u7a0b'],
   fs=15, c=C['text'])

rect(s, 7, 1.4, 5.5, 2.5, C['light'], C['accent'], 2, True)
tb(s, 7.2, 1.5, 5, 0.5, '\u9879\u76ee\u5b9a\u4f4d', fs=20, c=C['primary'], b=True, a=PP_ALIGN.CENTER)
mt(s, 7.3, 2.0, 5, 1.8,
   ['\u6a21\u5757\u4e00\uff1a\u8ba1\u7b97\u673a\u7f51\u7edc\uff08\u672c\u79d1\u00b7\u4e13\u4e1a\u6838\u5fc3\u8bfe\uff09',
    '\u5efa\u8bbe\u5468\u671f\uff1a2026.01 - 2027.12\uff082\u5e74\uff09',
    '\u6838\u5fc3\u6539\u9769\uff1a\u56fd\u4ea7\u4e91\u5e73\u53f0\u4e3a\u5e95\u5ea7\uff0c\u6df1\u5ea6\u878d\u5408\u667a\u80fd\u4f53\u5de5\u5177\uff0c',
    '\u6784\u5efa\u201c\u667a\u80fd\u751f\u6210-\u79c1\u6709\u7b54\u7591-\u667a\u6167\u7ba1\u7406-\u5b9e\u6218\u90e8\u7f72\u201d\u5168\u94fe\u6761\u6559\u5b66\u65b0\u8303\u5f0f'],
   fs=14, c=C['text'])

rect(s, 0.8, 4.3, 11.5, 2.8, C['primary'], r=True)
tb(s, 1, 4.4, 3, 0.5, '\u5b98\u65b9\u8fdb\u5ea6\u5b89\u6392', fs=18, c=C['accent'], b=True)
mt(s, 1, 4.95, 11, 2.0,
   ['\u2022 2026.1-12\uff1a\u8bfe\u7a0b\u5efa\u8bbe\u9636\u6bb5\uff08\u77e5\u8bc6\u4f53\u7cfb\u68b3\u7406\u3001\u6559\u6750\u64b0\u5199\u3001\u5b9e\u8df5\u8d44\u6e90\u5efa\u8bbe\uff09',
    '\u2022 2027.1-12\uff1a\u8bd5\u70b9\u4e0e\u5b8c\u5584\u9636\u6bb5\uff08\u4e0d\u5c11\u4e8e5\u6240\u9ad8\u6821\u8bd5\u7528\uff0c\u6839\u636e\u53cd\u9988\u8fed\u4ee3\uff09',
    '\u2022 2026.11-12\uff1a\u4e2d\u671f\u5de5\u4f5c\u603b\u7ed3\uff0c\u786e\u7acb\u201c\u56db\u4e2a\u6838\u5fc3\u201d\u5efa\u8bbe\u8d28\u91cf\u6807\u51c6\u4e0e\u8ba4\u5b9a\u529e\u6cd5'],
   fs=14, c=C['white'])

# ============================================================
# SLIDE 3: Team Foundation
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, C['white'])
hdr(s, '\u56e2\u961f\u73b0\u6709\u6210\u679c\u57fa\u7840')

ach = [
    { 't': 'MOOC\u5728\u7ebf\u8bfe\u7a0b',
      'd': '\u300a\u9ad8\u7ea7\u7f51\u7edc\u901a\u4fe1\u539f\u7406\u5b9e\u8df5\u300b\n\u4e2d\u56fd\u5927\u5b66MOOC | 228\u4eba\u5b66\u4e60\n18\u5468\u8bfe\u7a0b\uff08\u5df2\u5b8c\u621014\u5468\uff09\n\u5408\u4f5c\u4f01\u4e1a\uff1a\u963f\u91cc\u4e91', 'c': C['secondary'] },
    { 't': '\u5df2\u51fa\u7248\u6559\u6750',
      'd': '\u300a\u7f51\u7edc\u901a\u4fe1\u539f\u7406\u5b9e\u8df5\u300b\n\uff08\u6e05\u534e\u51fa\u7248\u793e,2024\uff09\n\u300a\u7cfb\u7edf\u7ea7\u7f16\u7a0b\u53ca\u5206\u5e03\u5f0f\u5e94\u7528\u5b9e\u73b0\u6280\u672f\u300b\n\uff08\u6e05\u534e\u51fa\u7248\u793e,\u5df2\u51fa\u7248\uff09', 'c': C['teal'] },
    { 't': '\u7701\u7ea7\u4e00\u6d41\u5b9e\u9a8c\u8bfe',
      'd': '\u300a\u8ba1\u7b97\u673a\u7f51\u7edc\u5b9e\u9a8c\u300b\n7\u671f\u8fd0\u884c | 623\u4eba\u6b21\u9009\u8bfe\n22\u6240\u9ad8\u6821\u8986\u76d6\n\u7d2f\u8ba1\u8bbf\u95ee59.4\u4e07\u6b21', 'c': C['success'] },
    { 't': '\u5b9e\u9a8c\u5e73\u53f0',
      'd': '\u963f\u91cc\u4e91\u4e91\u5b9e\u9a8c\u5ba4\u5e73\u53f0\nMininet / OpenDaylight / OpenStack\n\u652f\u6301\u7f51\u7edc\u865a\u62df\u5316\u5b9e\u6218\n\u4f01\u4e1a\u7ea7\u5b9e\u9a8c\u73af\u5883', 'c': C['accent'] }
]

ax = 0.5
for a in ach:
    rect(s, ax, 1.5, 3.05, 4.0, C['light'], a['c'], 2, True)
    tb(s, ax+0.15, 1.6, 2.75, 0.5, a['t'], fs=17, c=a['c'], b=True, a=PP_ALIGN.CENTER)
    mt(s, ax+0.15, 2.2, 2.75, 3.0, a['d'].split('\n'), fs=13, c=C['text'])
    ax += 3.2

tb(s, 0.8, 5.7, 11, 0.5,
   '\u56e2\u961f\u6838\u5fc3\u6210\u5458\uff1a\u8bf8\u845b\u658c\uff08\u6559\u6388\uff09\u3001\u91d1\u84c9\uff08\u526f\u6559\u6388\uff09\u3001\u9ad8\u660e\u3001\u674e\u4f20\u714c\u3001\u848b\u732e',
   fs=14, c=C['muted'], a=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 4: Core Reform Direction (AI Native Paradigm)
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, C['white'])
hdr(s, '\u6838\u5fc3\u6539\u9769\u65b9\u5411\uff1aAI\u539f\u751f\u6559\u5b66\u65b0\u8303\u5f0f')

tb(s, 0.8, 1.4, 3, 0.5, '\u7406\u5ff5\u8f6c\u53d8', fs=20, c=C['accent'], b=True)
mt(s, 0.8, 2.0, 5.5, 2.8,
   ['\u2022 \u4ece\u201c\u77e5\u8bc6\u704c\u8f93\u201d \u2192 \u201cAI\u5e94\u7528\u7684\u5b9e\u8df5\u573a\u666f\u201d',
    '\u2022 \u201c\u8bfe\u7a0b\u4e0d\u518d\u662f\u6559\u5b66\u6838\u5fc3\uff0c\u800c\u662fAI\u5de5\u5177\u5e94\u7528\u7684\u80cc\u666f\u573a\u666f\u201d',
    '\u2022 \u6838\u5fc3\u76ee\u6807\uff1a\u57f9\u517b\u5bf9AI\u751f\u6210\u5185\u5bb9\u7684\u5224\u65ad\u529b\u4e0e\u6279\u5224\u6027\u601d\u7ef4',
    '\u2022 \u4e94\u5e74\u540e\u4e13\u4e1a\u6559\u80b2\u5c06\u5168\u9762\u91cd\u6784',
    '\u2022 \u6559\u5e08\u9700\u5148\u76f8\u4fe1AI\u7684\u6f5c\u529b\uff0c\u624d\u80fd\u8c03\u6574\u6559\u5b66\u7b56\u7565'],
   fs=14, c=C['text'])

rect(s, 7, 1.3, 5.5, 3.5, C['primary'], r=True)
tb(s, 7.2, 1.4, 5, 0.5, '\u5168\u94fe\u6761\u6559\u5b66\u65b0\u8303\u5f0f', fs=20, c=C['accent'], b=True, a=PP_ALIGN.CENTER)
mt(s, 7.3, 2.1, 5, 2.5,
   ['\u2460 \u667a\u80fd\u751f\u6210 \u2014 AI\u8f85\u52a9\u77e5\u8bc6\u70b9\u5185\u5bb9\u751f\u6210',
    '\u2461 \u79c1\u6709\u7b54\u7591 \u2014 \u77e5\u8bc6\u5e93\u673a\u5668\u4eba24\u5c0f\u65f6\u7cbe\u51c6\u7b54\u7591',
    '\u2462 \u667a\u6167\u7ba1\u7406 \u2014 MOOC\u6155\u8bfe\u5802+\u9489\u9489AI\u52a9\u7406\u5168\u6d41\u7a0b\u8986\u76d6',
    '\u2463 \u5b9e\u6218\u90e8\u7f72 \u2014 \u56fd\u4ea7\u4e91\u5e73\u53f0\u771f\u5b9e\u73af\u5883\u9a8c\u8bc1'],
   fs=15, c=C['white'])

rect(s, 0.8, 5.2, 11.5, 1.3, C['light'], C['accent'], 2, True)
mt(s, 1, 5.25, 11, 1.2,
   ['\u201c\u672a\u6765\u8bfe\u7a0b\u5efa\u8bbe\u7684\u5173\u952e\u4e0d\u518d\u662f\u77e5\u8bc6\u4f20\u6388\uff0c\u800c\u662f\u57f9\u517b\u5b66\u751f\u5bf9AI\u751f\u6210\u5185\u5bb9\u7684\u5224\u65ad\u529b\u3002\u201d',
    '\u201c\u4ea4\u4e00\u767e\u4efd\u4f5c\u4e1a\uff0c\u53ea\u8981\u6709\u4e24\u4efd\u8d85\u51fa\u60f3\u8c61\uff0c\u5c31\u662f\u6210\u529f\u3002\u201d'],
   fs=15, c=C['primary'], italic=True, a=PP_ALIGN.CENTER)
tb(s, 1, 6.2, 11, 0.3, '\u2014\u2014 \u8bf8\u845b\u658c', fs=13, c=C['muted'], a=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 5: Five Reform Initiatives
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, C['white'])
hdr(s, '\u4e94\u5927\u6539\u9769\u4e3e\u63aa')

ini = [
    { 'n': '01', 't': '\u534f\u8bae\u673a\u5236\u53ef\u89c6\u5316',
      'd': '\u5229\u7528AI\u5de5\u5177\u751f\u6210\u52a8\u6001\u52a8\u753b\uff0c\u8986\u76d6STP\u3001\u8def\u7531\u67e5\u627e\u3001CHAP\u3001ACL\u3001NAT\u3001TCP\u4e09\u6b21\u63e1\u624b\u3001OpenFlow\u7b49\u6838\u5fc3\u534f\u8bae\u3002\u5b66\u751f\u4e3b\u5bfc\u751f\u6210\uff0c\u6559\u5e08\u7b5b\u9009\u6700\u4f18\uff0c\u5f62\u6210\u534f\u8bae\u52a8\u753b\u5e93\u3002', 'c': C['secondary'] },
    { 'n': '02', 't': '\u8bfe\u7a0b\u79c1\u6709\u77e5\u8bc6\u5e93\u673a\u5668\u4eba',
      'd': '\u6574\u5408\u6559\u6750\u3001\u8bfe\u4ef6\u3001\u4e60\u9898\u3001FAQ\uff0c\u63d0\u4f9b24\u5c0f\u65f6\u4e2a\u6027\u5316\u8f85\u5bfc\u7b54\u7591\u3002\u6280\u672f\u5b9e\u73b0\uff1a\u767e\u70bc\u77e5\u8bc6\u5e93\u00d7OpenClaw\u96c6\u6210\uff0c\u4e09\u7ea7\u8bb0\u5fc6\u67b6\u6784\uff08L1\u5de5\u4f5c/L2\u77e5\u8bc6/L3\u957f\u671f\uff09\u3002', 'c': C['teal'] },
    { 'n': '03', 't': '\u5168\u6d41\u7a0b\u667a\u6167\u6559\u5b66\u7ba1\u7406',
      'd': 'MOOC\u6155\u8bfe\u5802+\u9489\u9489AI\u52a9\u7406\uff0c\u8bfe\u524d-\u8bfe\u4e2d-\u8bfe\u540e\u5168\u6d41\u7a0b\u8986\u76d6\uff0c\u6570\u636e\u9a71\u52a8\u7684\u6559\u5b66\u6539\u8fdb\u3002\u901a\u8fc7\u9489\u9489\u7fa4\u4e0eAI\u52a9\u624b\u534f\u4f5c\uff0c\u5b9e\u73b0\u4f1a\u8bae\u7eaa\u8981\u81ea\u52a8\u751f\u6210\u3001\u8bfe\u7a0b\u603b\u7ed3\u667a\u80fd\u6c47\u603b\u3002', 'c': C['accent'] },
    { 'n': '04', 't': '\u667a\u80fd\u4f53\u8f85\u52a9\u5b9e\u9a8c\u65b0\u8303\u5f0f',
      'd': 'AI\u8f85\u52a9\u751f\u6210\u547d\u4ee4+\u4eba\u5de5\u4e13\u4e1a\u6392\u9519+\u771f\u5b9e\u8bbe\u5907\u90e8\u7f72\u3002\u4fdd\u7559\u771f\u5b9e\u4ea4\u4e92\u903b\u8f91\uff0cAI\u8f85\u52a9\u6982\u5ff5\u7406\u89e3\u3002\u9a8c\u8bc1\u6027/\u8bbe\u8ba1\u6027/\u7efc\u5408\u6027=3:4:3\uff0c\u5b9e\u9a8c\u6548\u7387\u63d0\u534740%\u3002', 'c': C['success'] },
    { 'n': '05', 't': '\u667a\u80fd\u4f53\u9a71\u52a8\u7684\u8bfe\u7a0b\u95e8\u6237\u5efa\u8bbe',
      'd': '\u8bfe\u7a0b\u5b98\u7f51\uff0c\u96c6\u6210\u6559\u5b66\u8d44\u6e90+\u5b9e\u9a8c\u6848\u4f8b+\u4f18\u79c0\u4f5c\u54c1+\u667a\u80fd\u4e92\u52a8\u3002\u5a92\u4f53\u5316\u8868\u8fbe\uff0c\u8d85\u8d8a\u4f20\u7edfPPT\u548c\u6559\u6750\u7684\u89c6\u89c9\u5448\u73b0\u9ad8\u5ea6\u3002', 'c': C['purple'] }
]

iy = 1.3
for it in ini:
    rect(s, 0.6, iy, 0.7, 1.0, it['c'], r=True)
    tb(s, 0.65, iy+0.15, 0.6, 0.7, it['n'], fs=20, c=C['white'], b=True, a=PP_ALIGN.CENTER)
    tb(s, 1.5, iy+0.05, 4.5, 0.4, it['t'], fs=16, c=it['c'], b=True)
    tb(s, 1.5, iy+0.45, 11.0, 0.55, it['d'], fs=12, c=C['text'])
    iy += 1.1

# ============================================================
# SLIDE 6: Student Assignment = Course Resource
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, C['white'])
hdr(s, '\u521b\u65b0\u6a21\u5f0f\uff1a\u201c\u5b66\u751f\u4f5c\u4e1a\u5373\u8bfe\u7a0b\u8d44\u6e90\u201d')

tb(s, 0.8, 1.4, 4, 0.5, '\u64cd\u4f5c\u6d41\u7a0b', fs=20, c=C['accent'], b=True)

steps = [
    ('\u2460 \u77e5\u8bc6\u70b9\u62c6\u89e3', '\u5c06\u8bfe\u7a0b\u77e5\u8bc6\u70b9\u62c6\u89e3\u4e3a\u5b66\u751f\u4f5c\u4e1a\u4efb\u52a1\uff0c\u6bcf\u8282\u8bfe1-2\u4e2a\u77e5\u8bc6\u70b9'),
    ('\u2461 \u5b66\u751fAI\u5b9e\u8df5', '\u5b66\u751f\u4f7f\u7528AI\u5de5\u5177\u751f\u6210\u53ef\u89c6\u5316\u5185\u5bb9\uff08\u52a8\u753b\u3001\u8bb2\u89e3\u3001\u4ee3\u7801\u7b49\uff09'),
    ('\u2462 \u6548\u679c\u7b5b\u9009', '\u6559\u5e08\u4ece\u5168\u73ed\u63d0\u4ea4\u7684AI\u751f\u6210\u5185\u5bb9\u4e2d\u6311\u9009\u6bcf\u77e5\u8bc6\u70b93\u4e2a\u6700\u4f18\u4f5c\u54c1'),
    ('\u2463 \u8d44\u6e90\u5e93\u5efa\u8bbe', '\u5f62\u6210\u53ef\u590d\u7528\u7684\u77e5\u8bc6\u70b9\u8d44\u6e90\u5e93\uff0c\u6bcf\u77e5\u8bc6\u70b9\u914d\u4e00\u4e2a\u4f18\u8d28AI\u751f\u6210\u5730\u5740')
]

sy = 2.0
for title, desc in steps:
    rect(s, 0.8, sy, 5.5, 0.9, C['light'], r=True)
    tb(s, 1.0, sy+0.05, 5.0, 0.35, title, fs=14, c=C['primary'], b=True)
    tb(s, 1.0, sy+0.4, 5.0, 0.45, desc, fs=11, c=C['text'])
    sy += 1.0

rect(s, 7, 1.3, 5.5, 3.8, C['light'], C['accent'], 2, True)
tb(s, 7.2, 1.4, 5, 0.5, '\u4ef7\u503c\u8bc4\u5224\u6743\u8f6c\u79fb', fs=18, c=C['primary'], b=True, a=PP_ALIGN.CENTER)
mt(s, 7.2, 2.0, 5.1, 3.0,
   ['\u5747\u503c\u5bfc\u5411 \u2192 \u5cf0\u503c\u53d1\u73b0',
    '  \u4e0d\u8981\u6c42\u6240\u6709\u4eba\u505a\u5bf9\uff0c\u800c\u662f\u53d1\u73b0\u4f18\u79c0\u4f5c\u54c1',
    '',
    '\u6559\u5e08\u4e3b\u89c2\u8bc4\u4ef7 \u2192 \u5b66\u751f\u4f5c\u54c1\u8d28\u91cf\u9a71\u52a8',
    '  \u4ee5\u4f5c\u54c1\u8d28\u91cf\u800c\u975e\u6559\u5e08\u504f\u597d\u4e3a\u6807\u51c6',
    '',
    '\u4e00\u6b21\u6027\u8bc4\u4ef7 \u2192 \u8fed\u4ee3\u4f18\u5316',
    '  \u4e0d\u8981\u6c42\u4e00\u6b21\u6027\u5b8c\u7f8e\uff0c\u901a\u8fc7\u6301\u7eed\u4f18\u5316\u63d0\u5347\u6548\u679c',
    '',
    '\u201c\u4e0d\u5728\u4e8eAI\u80fd\u5426\u5b8c\u6210\u4efb\u52a1\uff0c',
    '  \u800c\u5728\u4e8e\u82b1\u66f4\u591a\u65f6\u95f4\u540e\u662f\u5426\u80fd\u505a\u5f97\u66f4\u597d\u3002\u201d'],
   fs=13, c=C['text'])

rect(s, 0.8, 5.5, 11.5, 1.5, C['light'], C['secondary'], 1.5, True)
tb(s, 1, 5.55, 3, 0.4, 'SDP\u8bfe\u7a0b\u5b9e\u8df5\u6848\u4f8b', fs=16, c=C['secondary'], b=True)
mt(s, 1, 5.95, 11, 1.0,
   ['\u2022 \u77e5\u8bc6\u70b9\u62c6\u89e3\uff1a\u5c06SDP\u6838\u5fc3\u6982\u5ff5\uff08\u63a7\u5236\u5668\u3001\u6570\u636e\u5e73\u9762\u3001OpenFlow\u534f\u8bae\u7b49\uff09\u62c6\u89e3\u4e3a10\u4e2a\u4f5c\u4e1a\u4efb\u52a1',
    '\u2022 \u5b66\u751f\u5b9e\u8df5\uff1a\u4f7f\u7528AI\u5de5\u5177\u751f\u6210\u52a8\u753b\u3001\u8bb2\u89e3\u89c6\u9891\u3001\u4ee3\u7801\u5b9e\u73b0',
    '\u2022 \u6548\u679c\u7b5b\u9009\uff1a\u4ece100\u4efd\u4f5c\u4e1a\u4e2d\u7b5b\u9009\u51fa30\u4efd\u4f18\u8d28\u4f5c\u54c1\uff0c\u5f62\u6210SDP\u77e5\u8bc6\u70b9\u8d44\u6e90\u5e93'],
   fs=13, c=C['text'])

# ============================================================
# SLIDE 7: DingTalk AI-Native Management (NEW - with screenshots)
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, C['white'])
hdr(s, '\u9489\u9489\u7fa4AI\u539f\u751f\u8bfe\u7a0b\u5efa\u8bbe\u7ba1\u7406\u673a\u5236')

# Left: mechanism description
tb(s, 0.6, 1.3, 5.5, 0.5, '\u7ba1\u7406\u67b6\u6784', fs=20, c=C['primary'], b=True)
mt(s, 0.6, 1.85, 5.8, 4.5,
   ['\u25b6 \u57fa\u4e8e\u9489\u9489\u7fa4\u7684AI\u52a9\u624b\u534f\u4f5c\u6a21\u5f0f',
    '',
    '\u2022 AI\u52a9\u624b\uff1a\u201c\u5c0f\u9f99\u867e\u201d\uff08\u57fa\u4e8eOpenClaw+Alibaba AI\uff09',
    '\u2022 \u7fa4\u540d\u79f0\uff1a\u201c\u667a\u80fd\u4f53\u5c0f\u9f99\u867e\u6d4b\u8bd5\u201d',
    '\u2022 \u53c2\u4e0e\u8005\uff1a\u8bfe\u7a0b\u56e2\u961f\u5168\u4f53\u6210\u5458 + AI\u52a9\u624b',
    '',
    '\u25b6 \u6838\u5fc3\u80fd\u529b',
    '',
    '\u2022 \u4f1a\u8bae\u7eaa\u8981\u81ea\u52a8\u751f\u6210\u4e0e\u5f52\u6863',
    '\u2022 \u4efb\u52a1\u5206\u914d\u4e0e\u8fdb\u5ea6\u8ddf\u8e2a',
    '\u2022 \u6559\u5b66\u6750\u6599\u81ea\u52a8\u751f\u6210\u4e0e\u6253\u5305',
    '\u2022 \u8bfe\u7a0b\u6848\u4f8b\u81ea\u52a8\u5bfc\u51fa\u4e0e\u5f52\u6863',
    '\u2022 \u77e5\u8bc6\u5e93\u673a\u5668\u4eba24\u5c0f\u65f6\u5728\u7ebf\u7b54\u7591',
    '',
    '\u25b6 \u6548\u679c\u6570\u636e',
    '',
    '\u2022 \u5df2\u751f\u62107+\u4e2a\u5b8c\u6574\u6559\u5b66\u6848\u4f8b',
    '\u2022 \u7d2f\u8ba1\u5904\u7406200+\u6761\u6d88\u606f',
    '\u2022 \u751f\u6210156\u4e2a\u6587\u4ef6\uff08Word/PPT/ZIP\uff09',
    '\u2022 \u5e73\u5747\u54cd\u5e94\u65f6\u95f4 < 30\u79d2'],
   fs=13, c=C['text'])

# Right: screenshot placeholder area
rect(s, 6.8, 1.3, 5.8, 5.2, C['light'], C['secondary'], 2, True)
tb(s, 7, 1.4, 5.4, 0.4, '\u9489\u9489\u7fa4AI\u534f\u4f5c\u573a\u666f\u622a\u56fe', fs=16, c=C['secondary'], b=True, a=PP_ALIGN.CENTER)

# Add images if available
img_base = '/home/admin/.openclaw/workspace/诸葛斌与小龙虾完整案例库/teaching_cases'
# Try to add screenshot images
img1 = f'{img_base}/002_数据分析系统开发案例/03_输入材料/c1ea85c1-fb90-4c1e-b8dc-393232cd9888.png'
img2 = f'{img_base}/004_钉钉 AI 助理教材_完整版/02_案例参考/03_输入材料/模块 5_案例文档预览截图.png'

if os.path.exists(img1):
    s.shapes.add_picture(img1, Inches(7.0), Inches(2.0), Inches(5.4), Inches(2.0))

if os.path.exists(img2):
    s.shapes.add_picture(img2, Inches(7.0), Inches(4.2), Inches(5.4), Inches(2.0))

# ============================================================
# SLIDE 8: Three Core Tasks
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, C['white'])
hdr(s, '\u4e09\u5927\u6838\u5fc3\u4efb\u52a1\u4f53\u7cfb')

tasks = [
    { 'num': '\u4efb\u52a1\u4e00', 'title': '\u8bfe\u7a0b\u5efa\u8bbe', 'lead': '\u8bf8\u845b\u658c',
      'items': ['\u77e5\u8bc6\u4f53\u7cfb\u68b3\u7406\uff0850-60\u4e2a\u77e5\u8bc6\u70b9\uff09',
                '\u77e5\u8bc6\u70b9\u56fe\u8c31\u53ef\u89c6\u5316\u6784\u5efa',
                '\u8bfe\u7a0b\u5efa\u8bbe\u6307\u5357\u7f16\u5199',
                '\u672c\u5730\u5316\u6559\u5b66\u5927\u7eb2\u8bbe\u8ba1\uff08\u56fd\u4ea7\u4e91+\u667a\u80fd\u4f53\uff09',
                '\u6559\u5b66\u8d44\u6e90\u5305\u5f00\u53d1\uff08\u52a8\u753b/\u8bfe\u4ef6/\u4e60\u9898\uff09',
                '\u5b9e\u9a8c\u8bfe\u7a0b2.0\u5347\u7ea7',
                '\u8bfe\u7a0b\u95e8\u6237\u7f51\u7ad9\u5efa\u8bbe'],
      'c': C['secondary'] },
    { 'num': '\u4efb\u52a1\u4e8c', 'title': '\u6559\u6750\u7f16\u5199', 'lead': '\u8bf8\u845b\u658c+\u9ad8\u660e',
      'items': ['\u57fa\u4e8e\u73b0\u6709\u6559\u6750\u5347\u7ea7\uff08\u975e\u5168\u65b0\u7f16\u5199\uff09',
                '\u6570\u5b57\u6559\u6750\u914d\u5957\u89c6\u9891/\u52a8\u753b/\u4ea4\u4e92',
                '\u6bcf\u7ae0\u914d\u5957\u6848\u4f8b\u22653\u4e2a',
                '\u4f01\u4e1a\u771f\u5b9e\u6848\u4f8b\u5360\u6bd4\u226530%',
                '\u6e05\u534e\u5927\u5b66\u51fa\u7248\u793e\u5408\u4f5c\u7eed\u7b7e',
                '\u6837\u7ae0\u64b0\u5199 + \u6559\u6750\u521d\u7a3f',
                '\u6839\u636e\u8bd5\u7528\u53cd\u9988\u8fed\u4ee3\u4fee\u8ba2'],
      'c': C['teal'] },
    { 'num': '\u4efb\u52a1\u4e09', 'title': '\u5b9e\u8df5\u6848\u4f8b', 'lead': '\u8bf8\u845b\u658c+\u848b\u732e',
      'items': ['\u9a8c\u8bc1\u6027/\u8bbe\u8ba1\u6027/\u7efc\u5408\u6027=3:4:3',
                '\u667a\u80fd\u4f53\u8f85\u52a9\u5b9e\u9a8c\uff08AI\u751f\u6210\u547d\u4ee4+\u4eba\u5de5\u6392\u9519\uff09',
                '\u963f\u91cc\u4e91\u5e73\u53f0+\u591a\u4e91\u652f\u6301\u5347\u7ea7',
                '\u62d3\u5c55\u534e\u4e3a/\u534e\u4e09\u7b49\u4f01\u4e1a\u5408\u4f5c',
                '\u5b9e\u9a8c\u9879\u76ee\u6e05\u5355\uff08\u226510\u4e2a\uff09',
                '\u5b9e\u8df5\u6848\u4f8b\u7fa4\u5efa\u8bbe\uff08\u22655\u4e2a\uff09',
                '\u4ea7\u5b66\u7814\u5408\u4f5c\u8d44\u6e90\uff08\u4f01\u4e1a\u771f\u5b9e\u9879\u76ee\u22652\u4e2a\uff09'],
      'c': C['accent'] }
]

tx = 0.4
for t in tasks:
    rect(s, tx, 1.35, 4.1, 5.6, C['light'], t['c'], 2, True)
    tb(s, tx+0.15, 1.4, 3.8, 0.5, f"{t['num']}\uff1a{t['title']}", fs=18, c=t['c'], b=True, a=PP_ALIGN.CENTER)
    tb(s, tx+0.15, 1.9, 3.8, 0.4, f"\u8d1f\u8d23\u4eba\uff1a{t['lead']}", fs=13, c=C['muted'], a=PP_ALIGN.CENTER)
    mt(s, tx+0.15, 2.3, 3.8, 4.5, t['items'], fs=12, c=C['text'])
    tx += 4.25

# ============================================================
# SLIDE 9: Milestones
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, C['white'])
hdr(s, '\u5173\u952e\u91cc\u7a0b\u7891\u4e0e\u65f6\u95f4\u8282\u70b9')

ms = [
    { 'd': '2026.06.30', 'i': '\u77e5\u8bc6\u4f53\u7cfb\u521d\u7a3f+\u6559\u5b66\u5927\u7eb2\u672c\u5730\u5316+\u7b2c1\u6b21\u7814\u8ba8\u8bb0\u5f55', 'ok': True },
    { 'd': '2026.07.15', 'i': '\u6559\u6750\u6846\u67b6\u8bbe\u8ba1+\u77e5\u8bc6\u70b9\u56fe\u8c31\u53ef\u89c6\u5316\u5b8c\u6210', 'ok': False },
    { 'd': '2026.07.31', 'i': '\u6559\u6750\u6837\u7ae02\u7ae0+\u5b9e\u9a8c\u6848\u4f8b5\u4e2a+\u5b9e\u9a8c\u9879\u76ee\u6e05\u5355', 'ok': False },
    { 'd': '2026.08.31', 'i': '\u534f\u8bae\u52a8\u753b5\u4e2a+\u77e5\u8bc6\u5e93\u673a\u5668\u4ebaMVP+\u4f01\u4e1a\u5408\u4f5c\u610f\u5411', 'ok': False },
    { 'd': '2026.09.30', 'i': '\u5168\u90e8\u4ea4\u4ed8\u7269\u521d\u7a3f+\u6848\u4f8b\u8d44\u6e90\u5e93+\u76ee\u6807\u9ad8\u6821\u8054\u7edc', 'ok': False },
    { 'd': '2026.10.31', 'i': '\u6559\u6750\u521d\u7a3f+\u5b9e\u9a8c\u8bfe\u7a0b2.0+\u4ea7\u5b66\u7814\u8d44\u6e90\u5e93+\u4e2d\u671f\u6750\u6599', 'ok': False },
    { 'd': '2026.11-12', 'i': '\u4e2d\u671f\u5de5\u4f5c\u603b\u7ed3+\u201c\u56db\u4e2a\u6838\u5fc3\u201d\u8d28\u91cf\u6807\u51c6\u5236\u5b9a', 'ok': False },
]

my = 1.4
for m in ms:
    dc = C['success'] if m['ok'] else C['light']
    dl = C['success'] if m['ok'] else C['muted']
    dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1), Inches(my), Inches(0.35), Inches(0.35))
    dot.fill.solid(); dot.fill.fore_color.rgb = dc; dot.line.color.rgb = dl; dot.line.width = Pt(1.5)
    tb(s, 1.5, my-0.02, 2.2, 0.4, m['d'], fs=16, c=C['primary'], b=True)
    tb(s, 3.8, my-0.02, 8, 0.4, m['i'], fs=14, c=C['text'])
    st = '\u2713 \u5df2\u5b8c\u6210' if m['ok'] else '\u25cb \u63a8\u8fdb\u4e2d'
    sc = C['success'] if m['ok'] else C['muted']
    tb(s, 12, my-0.02, 1.2, 0.4, st, fs=12, c=sc, b=True)
    my += 0.8

# ============================================================
# SLIDE 10: Wisdom Tree + Collaboration
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, C['white'])
hdr(s, '\u667a\u6167\u6811\u90e8\u7f72\u4e0e\u534f\u4f5c\u673a\u5236')

tb(s, 0.8, 1.4, 3, 0.5, '\u667a\u6167\u6811\u90e8\u7f72\u8fdb\u5ea6', fs=18, c=C['primary'], b=True)
dl = [
    ('\u2713 \u8bfe\u7a0b\u6846\u67b6\u642d\u5efa', C['success']),
    ('\u2713 \u6559\u5b66\u89c6\u9891\u4e0a\u4f20', C['success']),
    ('\u2713 \u7ae0\u8282\u6d4b\u9a8c\u914d\u7f6e', C['success']),
    ('\u27f3 \u8ba8\u8bba\u533a\u4e0e\u4e92\u52a8\u6a21\u5757\u8bbe\u7f6e', C['accent']),
    ('\u27f3 \u4f5c\u4e1a\u6279\u6539\u89c4\u5219\u914d\u7f6e', C['accent']),
    ('\u25cb AI\u8f85\u52a9\u7b54\u7591\u6a21\u5757\u63a5\u5165', C['muted']),
    ('\u25cb \u671f\u672b\u8003\u8bd5\u4e0e\u8bc4\u4ef7\u914d\u7f6e', C['muted']),
]
dy = 2.0
for txt, col in dl:
    tb(s, 0.8, dy, 5.5, 0.45, txt, fs=14, c=col, b='\u2713' in txt)
    dy += 0.45

tb(s, 7, 1.4, 4, 0.5, '\u534f\u4f5c\u673a\u5236', fs=18, c=C['primary'], b=True)
mt(s, 7, 2.0, 5.5, 3.0,
   ['\u2022 \u7ebf\u4e0a\u78b0\u5934\u4f1a\uff1a\u53cc\u5468 | \u5468\u4e94 15:00-16:00\uff08\u8bf8\u845b\u658c\uff09',
    '\u2022 \u7ebf\u4e0b\u7814\u8ba8\u4f1a\uff1a\u6bcf\u5b63\u5ea6\uff08\u5168\u4f53\uff09',
    '\u2022 \u6587\u6863\u5f52\u6863\uff1a\u6bcf\u6708\u5e95\uff08\u91d1\u84c9\uff09',
    '\u2022 \u4e2d\u671f\u603b\u7ed3\uff1a2026\u5e7411-12\u6708\uff08\u674e\u4f20\u714c\uff09',
    '',
    '\u2022 \u6587\u6863\u7edf\u4e00\u5b58\u653e\u9489\u9489\u7fa4\u6587\u4ef6/\u4e91\u5e73\u53f0',
    '\u2022 \u6bcf\u6708\u5e95\u63d0\u4ea4\u8fdb\u5ea6\u62a5\u544a',
    '\u2022 AI\u7b97\u529b\u5145\u88d5\u652f\u6301\u957f\u671f\u63a2\u7d22'],
   fs=14, c=C['text'])

rect(s, 0.8, 5.2, 11.5, 1.5, C['light'], C['accent'], 2, True)
tb(s, 1, 5.25, 3, 0.4, '\u8bd5\u70b9\u9ad8\u6821\u62d3\u5c55', fs=16, c=C['primary'], b=True)
mt(s, 1, 5.65, 11, 1.0,
   ['\u2022 \u6d59\u6c5f\u5de5\u4e1a\u5927\u5b66 | \u676d\u5dde\u7535\u5b50\u79d1\u6280\u5927\u5b66 | \u6d59\u6c5f\u7406\u5de5\u5927\u5b66 | \u5b81\u6ce2\u5927\u5b66 | \u6d59\u6c5f\u5e08\u8303\u5927\u5b66',
    '\u2022 \u8bd5\u7528\u5185\u5bb9\uff1a\u6559\u6750+\u5b9e\u9a8c+\u5e73\u53f0+\u667a\u80fd\u4f53\u7b54\u7591 | \u76ee\u6807\u6837\u672c\uff1a\u4e0d\u5c11\u4e8e500\u4efd\u53cd\u9988'],
   fs=13, c=C['text'])

# ============================================================
# SLIDE 11: Expected Outcomes
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, C['white'])
hdr(s, '\u9884\u671f\u6210\u679c\u4e0e\u5c55\u671b')

outcomes = [
    '\u5efa\u6210\u7b26\u5408\u201c101\u8ba1\u5212\u201d\u6807\u51c6\u7684\u8ba1\u7b97\u673a\u7f51\u7edc\u6838\u5fc3\u8bfe\u7a0b\u4f53\u7cfb',
    '\u51fa\u7248\u914d\u5957\u6570\u667a\u5316\u6559\u6750\uff08\u6e05\u534e\u5927\u5b66\u51fa\u7248\u793e\uff09\u4e0e\u534f\u8bae\u52a8\u753b\u5e93',
    '\u5f62\u6210\u8bfe\u7a0b\u79c1\u6709\u77e5\u8bc6\u5e93\u673a\u5668\u4eba + \u5168\u6d41\u7a0b\u667a\u6167\u6559\u5b66\u7ba1\u7406\u4f53\u7cfb',
    '\u5b8c\u6210\u667a\u6167\u6811\u5e73\u53f0\u5728\u7ebf\u8bfe\u7a0b\u90e8\u7f72\u4e0e\u8de8\u6821\u8fd0\u8425',
    '\u5b9e\u9a8c\u6548\u7387\u63d0\u534740%+\uff0c\u6bcf\u5b66\u671f\u53d7\u76ca200+\u5b66\u751f',
    '\u4e0d\u5c11\u4e8e5\u6240\u9ad8\u6821\u8bd5\u70b9\uff0c\u8986\u76d6500+\u5b66\u751f\u53cd\u9988\u6570\u636e',
    '\u6253\u9020\u201c101\u8ba1\u5212\u201d\u672c\u5730\u5316\u201cAI+\u6559\u80b2\u201d\u53ef\u590d\u5236\u3001\u53ef\u63a8\u5e7f\u5178\u578b\u6848\u4f8b',
    '\u57fa\u4e8e\u9489\u9489\u7fa4\u7684AI\u539f\u751f\u8bfe\u7a0b\u7ba1\u7406\u673a\u5236\u6210\u4e3a\u6559\u80b2\u884c\u4e1a\u65b0\u6807\u6746',
]

oy = 1.3
for o in outcomes:
    rect(s, 0.8, oy, 0.55, 0.55, C['accent'], r=True)
    tb(s, 1.55, oy-0.03, 10.5, 0.6, o, fs=15, c=C['text'], b=True)
    oy += 0.75

# ============================================================
# SLIDE 12: Thank You
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, C['primary'])
rect(s, 0, 0, 13.33, 0.12, C['accent'])
tb(s, 1, 2.0, 11.33, 1.0,
   '\u611f\u8c22\u8046\u542c\uff01',
   fs=44, c=C['white'], b=True, a=PP_ALIGN.CENTER)
tb(s, 1, 3.0, 11.33, 0.7,
   '\u6b22\u8fce\u6279\u8bc4\u6307\u6b63',
   fs=24, c=C['accent'], a=PP_ALIGN.CENTER)
mt(s, 1, 5.5, 11.33, 0.8,
   ['\u8bf8\u845b\u658c\u56e2\u961f \u00b7 \u6d59\u6c5f\u5de5\u5546\u5927\u5b66',
    '\u8ba1\u7b97\u673a\u7f51\u7edc\u201c101\u8ba1\u5212\u201d\u8bfe\u7a0b\u5efa\u8bbe\u9879\u76ee'],
   fs=16, c=C['muted'], a=PP_ALIGN.CENTER)

out = '/home/admin/.openclaw/workspace/101-network-course-v4-dingtalk.pptx'
prs.save(out)
print(f'PPT saved: {out}')
