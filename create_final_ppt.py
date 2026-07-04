#!/usr/bin/env python3
"""
生成购买力平价检验分析 PPT 演示文稿
包含6张图表：折线图、柱形图等
数据：2026年6月2日最新汇率
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import os

prs = Presentation()
prs.slide_width = Emu(12192000)  # 16:9
prs.slide_height = Emu(6858000)

CHART_DIR = '/home/admin/.openclaw/workspace/ppp_charts'

# ============================================================
# 配色方案
# ============================================================
BLUE = RGBColor(0x21, 0x96, 0xF3)
ORANGE = RGBColor(0xFF, 0x57, 0x22)
GREEN = RGBColor(0x4C, 0xAF, 0x50)
RED = RGBColor(0xF4, 0x43, 0x36)
DARK_BLUE = RGBColor(0x1A, 0x23, 0x7E)
GRAY = RGBColor(0x60, 0x60, 0x60)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT_BLUE = RGBColor(0x0D, 0x47, 0xA1)

# ============================================================
# 辅助函数
# ============================================================

def add_background(slide, color=WHITE):
    """设置幻灯片背景"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape_bg(slide, left, top, width, height, color):
    """添加形状背景"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18, 
                 color=DARK_BLUE, bold=False, alignment=PP_ALIGN.LEFT, font_name='微软雅黑'):
    """添加文本框"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_multiline_text(slide, left, top, width, height, lines, font_size=14, 
                       color=GRAY, line_spacing=1.5, bullet=False):
    """添加多行文本"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = '微软雅黑'
        p.space_after = Pt(font_size * line_spacing * 0.5)
        if bullet and i > 0:
            p.text = '• ' + line
    return txBox

def add_chart_image(slide, image_path, left, top, width, height):
    """添加图表图片"""
    if os.path.exists(image_path):
        slide.shapes.add_picture(image_path, left, top, width, height)
        return True
    return False

def add_accent_bar(slide, left, top, width, height, color=BLUE):
    """添加装饰条"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

# ============================================================
# 第1页：封面
# ============================================================
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
add_background(slide1, WHITE)

# 顶部装饰条
add_accent_bar(slide1, Emu(0), Emu(0), Emu(12192000), Emu(300000), ACCENT_BLUE)

# 主标题区域
add_shape_bg(slide1, Emu(0), Emu(2000000), Emu(12192000), Emu(3000000), LIGHT_GRAY)

add_text_box(slide1, Emu(1000000), Emu(2200000), Emu(10192000), Emu(800000),
             '购买力平价检验分析', font_size=44, color=DARK_BLUE, bold=True)
add_text_box(slide1, Emu(1000000), Emu(3000000), Emu(10192000), Emu(600000),
             '——德美日英实际汇率分析', font_size=28, color=BLUE, bold=False)

# 数据说明
add_text_box(slide1, Emu(1000000), Emu(4200000), Emu(10192000), Emu(400000),
             f'📊 数据日期：2026年6月2日（实时汇率）', font_size=16, color=GRAY)
add_text_box(slide1, Emu(1000000), Emu(4600000), Emu(10192000), Emu(400000),
             '📈 包含6张数据图表 | 名义汇率 vs PPP汇率 | 实际汇率检验', font_size=14, color=GRAY)

# 底部装饰
add_accent_bar(slide1, Emu(0), Emu(6558000), Emu(12192000), Emu(300000), ACCENT_BLUE)

# ============================================================
# 第2页：核心发现（摘要）
# ============================================================
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide2, WHITE)
add_accent_bar(slide2, Emu(0), Emu(0), Emu(12192000), Emu(150000), ACCENT_BLUE)

add_text_box(slide2, Emu(500000), Emu(300000), Emu(11192000), Emu(600000),
             '📌 核心发现', font_size=32, color=DARK_BLUE, bold=True)

# 四个关键发现卡片
cards = [
    ('🔴 日元偏离最大', '+42.5% 高估', '名义汇率 159.6 vs PPP 112.0', RED),
    ('🟡 欧元显著偏离', '+1.2% 高估', '名义汇率 0.86 vs PPP 0.85', ORANGE),
    ('🟢 英镑基本均衡', '-2.2% 低估', '名义汇率 0.743 vs PPP 0.76', GREEN),
    ('🟣 瑞郎明显高估', '+9.2% 高估', '名义汇率 0.786 vs PPP 0.72', RGBColor(0x9C, 0x27, 0xB0)),
]

for i, (title, subtitle, detail, color) in enumerate(cards):
    x = Emu(500000 + i * 2800000)
    y = Emu(1100000)
    
    # 卡片背景
    card = add_shape_bg(slide2, x, y, Emu(2600000), Emu(2200000), LIGHT_GRAY)
    
    # 顶部颜色条
    add_accent_bar(slide2, x, y, Emu(2600000), Emu(100000), color)
    
    # 内容
    add_text_box(slide2, x + Emu(100000), y + Emu(200000), Emu(2400000), Emu(500000),
                 title, font_size=20, color=color, bold=True)
    add_text_box(slide2, x + Emu(100000), y + Emu(700000), Emu(2400000), Emu(500000),
                 subtitle, font_size=28, color=DARK_BLUE, bold=True)
    add_text_box(slide2, x + Emu(100000), y + Emu(1300000), Emu(2400000), Emu(400000),
                 detail, font_size=14, color=GRAY)

# 结论框
add_shape_bg(slide2, Emu(500000), Emu(3500000), Emu(11192000), Emu(1500000), RGBColor(0xE3, 0xF2, 0xFD))
add_text_box(slide2, Emu(700000), Emu(3600000), Emu(10792000), Emu(400000),
             '⚠️ 结论：绝对购买力平价在短期不成立', font_size=22, color=DARK_BLUE, bold=True)
add_text_box(slide2, Emu(700000), Emu(4000000), Emu(10792000), Emu(800000),
             '所有货币名义汇率与PPP汇率均存在偏离，其中日元偏离最严重（+42.5%），\n'
             '英镑最接近均衡（-2.2%）。实际汇率持续偏离1.0，PPP在中长期也不完全成立。',
             font_size=16, color=GRAY)

# ============================================================
# 第3页：名义汇率 vs PPP汇率对比（图表1）
# ============================================================
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide3, WHITE)
add_accent_bar(slide3, Emu(0), Emu(0), Emu(12192000), Emu(150000), ACCENT_BLUE)

add_text_box(slide3, Emu(500000), Emu(300000), Emu(11192000), Emu(600000),
             '📊 名义汇率 vs PPP汇率对比', font_size=28, color=DARK_BLUE, bold=True)

# 插入图表
add_chart_image(slide3, f'{CHART_DIR}/chart1_nominal_vs_ppp.png',
                Emu(800000), Emu(1000000), Emu(10592000), Emu(5000000))

# 底部说明
add_text_box(slide3, Emu(500000), Emu(6200000), Emu(11192000), Emu(400000),
             '数据来源：exchangerate-api.com (2026.06.02) | 世界银行ICP | 柱形图对比显示名义汇率与PPP汇率差异',
             font_size=12, color=GRAY)

# ============================================================
# 第4页：PPP偏离度分析（图表2）
# ============================================================
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide4, WHITE)
add_accent_bar(slide4, Emu(0), Emu(0), Emu(12192000), Emu(150000), ACCENT_BLUE)

add_text_box(slide4, Emu(500000), Emu(300000), Emu(11192000), Emu(600000),
             '📉 PPP偏离度分析', font_size=28, color=DARK_BLUE, bold=True)

# 插入图表
add_chart_image(slide4, f'{CHART_DIR}/chart2_ppp_deviation.png',
                Emu(800000), Emu(1000000), Emu(10592000), Emu(5000000))

# 底部说明
add_text_box(slide4, Emu(500000), Emu(6200000), Emu(11192000), Emu(400000),
             '绿色=高估 | 红色=低估 | 偏离度 = (名义汇率-PPP汇率)/PPP汇率 × 100%',
             font_size=12, color=GRAY)

# ============================================================
# 第5页：历史汇率走势（图表3）
# ============================================================
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide5, WHITE)
add_accent_bar(slide5, Emu(0), Emu(0), Emu(12192000), Emu(150000), ACCENT_BLUE)

add_text_box(slide5, Emu(500000), Emu(300000), Emu(11192000), Emu(600000),
             '📈 历史汇率走势 (2020-2026)', font_size=28, color=DARK_BLUE, bold=True)

# 插入图表
add_chart_image(slide5, f'{CHART_DIR}/chart3_historical_rates.png',
                Emu(500000), Emu(1000000), Emu(11192000), Emu(5200000))

# 底部说明
add_text_box(slide5, Emu(500000), Emu(6300000), Emu(11192000), Emu(400000),
             '折线图展示2020-2026年主要货币兑美元汇率变化 | 数据来源：各国央行年度平均汇率',
             font_size=12, color=GRAY)

# ============================================================
# 第6页：实际汇率检验（图表4）
# ============================================================
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide6, WHITE)
add_accent_bar(slide6, Emu(0), Emu(0), Emu(12192000), Emu(150000), ACCENT_BLUE)

add_text_box(slide6, Emu(500000), Emu(300000), Emu(11192000), Emu(600000),
             '⚖️ 实际汇率与购买力平价检验', font_size=28, color=DARK_BLUE, bold=True)

# 插入图表
add_chart_image(slide6, f'{CHART_DIR}/chart4_real_exchange_rate.png',
                Emu(800000), Emu(1000000), Emu(10592000), Emu(5000000))

# 底部说明
add_text_box(slide6, Emu(500000), Emu(6200000), Emu(11192000), Emu(400000),
             '实际汇率 q = 名义汇率/PPP汇率 | q>1.05高估 | q<0.95低估 | q≈1均衡',
             font_size=12, color=GRAY)

# ============================================================
# 第7页：价格水平对比（图表5）
# ============================================================
slide7 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide7, WHITE)
add_accent_bar(slide7, Emu(0), Emu(0), Emu(12192000), Emu(150000), ACCENT_BLUE)

add_text_box(slide7, Emu(500000), Emu(300000), Emu(11192000), Emu(600000),
             '💰 各国价格水平对比', font_size=28, color=DARK_BLUE, bold=True)

# 插入图表
add_chart_image(slide7, f'{CHART_DIR}/chart5_cpi_comparison.png',
                Emu(800000), Emu(1000000), Emu(10592000), Emu(5000000))

# 底部说明
add_text_box(slide7, Emu(500000), Emu(6200000), Emu(11192000), Emu(400000),
             'CPI指数（美国=100）| 瑞士价格水平最高，日本最低 | 数据来源：各国统计局 2025年',
             font_size=12, color=GRAY)

# ============================================================
# 第8页：归一化趋势（图表6）
# ============================================================
slide8 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide8, WHITE)
add_accent_bar(slide8, Emu(0), Emu(0), Emu(12192000), Emu(150000), ACCENT_BLUE)

add_text_box(slide8, Emu(500000), Emu(300000), Emu(11192000), Emu(600000),
             '📊 汇率变化趋势（归一化对比）', font_size=28, color=DARK_BLUE, bold=True)

# 插入图表
add_chart_image(slide8, f'{CHART_DIR}/chart6_normalized_trend.png',
                Emu(500000), Emu(1000000), Emu(11192000), Emu(5200000))

# 底部说明
add_text_box(slide8, Emu(500000), Emu(6300000), Emu(11192000), Emu(400000),
             '归一化至2020年=100 | 折线图清晰展示各货币相对强弱变化趋势',
             font_size=12, color=GRAY)

# ============================================================
# 第9页：理论检验结果
# ============================================================
slide9 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide9, WHITE)
add_accent_bar(slide9, Emu(0), Emu(0), Emu(12192000), Emu(150000), ACCENT_BLUE)

add_text_box(slide9, Emu(500000), Emu(300000), Emu(11192000), Emu(600000),
             '📋 购买力平价理论检验结果', font_size=28, color=DARK_BLUE, bold=True)

# 检验结果表格
tests = [
    ('绝对PPP检验', '名义汇率 vs PPP汇率', '❌ 不成立', '所有货币均存在显著偏离'),
    ('实际汇率检验', '实际汇率 vs 1.0', '❌ 不成立', '实际汇率持续偏离均衡值'),
    ('相对PPP检验', '汇率变化率 vs 通胀差', '⚠️ 部分成立', '中长期有一定解释力'),
]

for i, (test, method, result, detail) in enumerate(tests):
    y = Emu(1100000 + i * 1500000)
    
    # 行背景
    bg_color = LIGHT_GRAY if i % 2 == 0 else WHITE
    add_shape_bg(slide9, Emu(500000), y, Emu(11192000), Emu(1200000), bg_color)
    
    add_text_box(slide9, Emu(600000), y + Emu(100000), Emu(2500000), Emu(400000),
                 test, font_size=20, color=DARK_BLUE, bold=True)
    add_text_box(slide9, Emu(3200000), y + Emu(100000), Emu(3000000), Emu(400000),
                 method, font_size=16, color=GRAY)
    add_text_box(slide9, Emu(6300000), y + Emu(100000), Emu(1500000), Emu(400000),
                 result, font_size=18, color=RED if '❌' in result else ORANGE, bold=True)
    add_text_box(slide9, Emu(7900000), y + Emu(100000), Emu(3500000), Emu(400000),
                 detail, font_size=14, color=GRAY)

# 总结
add_shape_bg(slide9, Emu(500000), Emu(5600000), Emu(11192000), Emu(800000), RGBColor(0xFF, 0xF3, 0xE0))
add_text_box(slide9, Emu(700000), Emu(5650000), Emu(10792000), Emu(700000),
             '💡 结论：绝对购买力平价在短期和中期均不成立，实际汇率呈现趋势性偏离。\n'
             '   相对PPP在中长期具有一定解释力，但需结合其他因素综合分析。',
             font_size=16, color=DARK_BLUE)

# ============================================================
# 第10页：偏离原因分析
# ============================================================
slide10 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide10, WHITE)
add_accent_bar(slide10, Emu(0), Emu(0), Emu(12192000), Emu(150000), ACCENT_BLUE)

add_text_box(slide10, Emu(500000), Emu(300000), Emu(11192000), Emu(600000),
             '🔍 PPP偏离原因分析', font_size=28, color=DARK_BLUE, bold=True)

causes = [
    ('1', '巴拉萨-萨缪尔森效应', '生产率差异导致可贸易品与不可贸易品价格差异，\n高生产率增长国家实际汇率倾向于升值'),
    ('2', '非贸易品价格差异', '服务业、房地产等非贸易品价格不受套利机制约束，\n各国非贸易品价格差异显著'),
    ('3', '贸易成本与壁垒', '运输成本、关税、非关税壁垒阻碍一价定律实现，\n贸易品价格仍存在系统性差异'),
    ('4', '资本流动与投机', '短期资本流动对汇率的影响超过贸易流量，\n市场预期和投机行为导致汇率超调'),
    ('5', '市场不完全竞争', '价格粘性导致价格调整滞后，\n企业定价策略导致同种商品价格差异'),
]

for i, (num, title, desc) in enumerate(causes):
    y = Emu(1000000 + i * 1100000)
    
    # 数字圆圈
    circle = slide10.shapes.add_shape(MSO_SHAPE.OVAL, Emu(600000), y + Emu(50000), Emu(600000), Emu(600000))
    circle.fill.solid()
    circle.fill.fore_color.rgb = BLUE
    circle.line.fill.background()
    tf = circle.text_frame
    tf.paragraphs[0].text = num
    tf.paragraphs[0].font.size = Pt(24)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # 标题
    add_text_box(slide10, Emu(1400000), y, Emu(4000000), Emu(400000),
                 title, font_size=20, color=DARK_BLUE, bold=True)
    # 描述
    add_text_box(slide10, Emu(1400000), y + Emu(400000), Emu(10000000), Emu(600000),
                 desc, font_size=14, color=GRAY)

# ============================================================
# 第11页：政策含义
# ============================================================
slide11 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide11, WHITE)
add_accent_bar(slide11, Emu(0), Emu(0), Emu(12192000), Emu(150000), ACCENT_BLUE)

add_text_box(slide11, Emu(500000), Emu(300000), Emu(11192000), Emu(600000),
             '🏛️ 政策含义与建议', font_size=28, color=DARK_BLUE, bold=True)

policies = [
    ('汇率政策制定', '不能仅依赖PPP作为汇率锚定目标，需综合考虑经济基本面、\n资本流动和市场预期。各国应根据自身情况选择合适的汇率制度。'),
    ('国际比较', '使用PPP汇率进行GDP等国际比较时，需注意PPP汇率与名义汇率的\n系统性差异。建议同时报告两种汇率下的比较结果。'),
    ('投资决策', '实际汇率偏离为跨境投资提供参考，但需结合其他基本面因素。\n高估货币国家的资产可能存在回调风险。'),
    ('长期趋势', '尽管短期偏离显著，长期来看PPP仍具有一定的均值回归特征。\n建议投资者关注实际汇率的长期趋势变化。'),
]

for i, (title, desc) in enumerate(policies):
    y = Emu(1000000 + i * 1400000)
    
    # 卡片背景
    card = add_shape_bg(slide11, Emu(500000), y, Emu(11192000), Emu(1200000), LIGHT_GRAY)
    
    # 标题
    add_text_box(slide11, Emu(700000), y + Emu(100000), Emu(10792000), Emu(400000),
                 f'📌 {title}', font_size=20, color=DARK_BLUE, bold=True)
    # 描述
    add_text_box(slide11, Emu(700000), y + Emu(500000), Emu(10792000), Emu(600000),
                 desc, font_size=15, color=GRAY)

# ============================================================
# 第12页：结论
# ============================================================
slide12 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide12, WHITE)
add_accent_bar(slide12, Emu(0), Emu(0), Emu(12192000), Emu(150000), ACCENT_BLUE)

add_text_box(slide12, Emu(500000), Emu(300000), Emu(11192000), Emu(600000),
             '✅ 研究结论', font_size=28, color=DARK_BLUE, bold=True)

conclusions = [
    ('1', '绝对购买力平价不成立', '名义汇率与PPP汇率存在系统性偏离，日本偏离最大（+42.5%），\n英国最小（-2.2%）。短期和中期均不成立。'),
    ('2', '实际汇率呈现趋势性', '实际汇率不是围绕1.0随机波动，存在持续的高估或低估现象。\n巴拉萨-萨缪尔森效应是重要解释因素。'),
    ('3', '相对PPP更有解释力', '相对PPP（汇率变化率≈通胀差）在中长期具有一定解释力，\n但需结合资本流动、市场预期等因素综合分析。'),
    ('4', '政策应用需谨慎', 'PPP可作为长期参考，但不宜作为短期政策依据。建议结合\n多种指标综合判断汇率均衡水平。'),
]

for i, (num, title, desc) in enumerate(conclusions):
    y = Emu(1000000 + i * 1400000)
    
    bg_color = RGBColor(0xE8, 0xF5, 0xE9) if i % 2 == 0 else WHITE
    add_shape_bg(slide12, Emu(500000), y, Emu(11192000), Emu(1200000), bg_color)
    
    # 数字
    circle = slide12.shapes.add_shape(MSO_SHAPE.OVAL, Emu(600000), y + Emu(100000), Emu(500000), Emu(500000))
    circle.fill.solid()
    circle.fill.fore_color.rgb = GREEN
    circle.line.fill.background()
    tf = circle.text_frame
    tf.paragraphs[0].text = num
    tf.paragraphs[0].font.size = Pt(22)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    add_text_box(slide12, Emu(1300000), y + Emu(100000), Emu(10000000), Emu(400000),
                 title, font_size=20, color=DARK_BLUE, bold=True)
    add_text_box(slide12, Emu(1300000), y + Emu(500000), Emu(10000000), Emu(600000),
                 desc, font_size=15, color=GRAY)

# ============================================================
# 第13页：数据来源与AI提示词
# ============================================================
slide13 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide13, WHITE)
add_accent_bar(slide13, Emu(0), Emu(0), Emu(12192000), Emu(150000), ACCENT_BLUE)

add_text_box(slide13, Emu(500000), Emu(300000), Emu(11192000), Emu(600000),
             '📚 数据来源与AI提示词', font_size=28, color=DARK_BLUE, bold=True)

# 数据来源
add_text_box(slide13, Emu(500000), Emu(1000000), Emu(5500000), Emu(400000),
             '📊 数据来源', font_size=22, color=DARK_BLUE, bold=True)

sources = [
    '• 名义汇率：exchangerate-api.com (2026.06.02)',
    '• PPP汇率：世界银行国际比较项目(ICP)',
    '• CPI指数：各国统计局、OECD (2025年)',
    '• 历史汇率：各国央行年度数据',
]

add_multiline_text(slide13, Emu(500000), Emu(1400000), Emu(5500000), Emu(2000000),
                   sources, font_size=15, color=GRAY)

# AI提示词
add_text_box(slide13, Emu(6200000), Emu(1000000), Emu(5500000), Emu(400000),
             '🤖 AI提示词（供进一步分析）', font_size=22, color=DARK_BLUE, bold=True)

prompts = [
    '1. 时间序列分析：2015-2024年PPP检验',
    '2. 巴拉萨-萨缪尔森效应检验',
    '3. 基于PPP的汇率预测模型',
    '4. 政策影响分析',
    '5. 一篮子货币PPP比较',
]

add_multiline_text(slide13, Emu(6200000), Emu(1400000), Emu(5500000), Emu(2000000),
                   prompts, font_size=15, color=GRAY)

# 参考文献
add_shape_bg(slide13, Emu(500000), Emu(3800000), Emu(11192000), Emu(2800000), LIGHT_GRAY)
add_text_box(slide13, Emu(700000), Emu(3900000), Emu(10792000), Emu(400000),
             '📖 参考文献', font_size=20, color=DARK_BLUE, bold=True)

refs = [
    '• Cassel, G. (1918). The Present Situation of the Foreign Exchanges.',
    '• Balassa, B. (1964). The Purchasing-Power-Parity Theory.',
    '• Rogoff, K. (1996). The Purchasing Power Parity Puzzle.',
    '• 世界银行. (2021). International Comparison Program (ICP).',
    '• IMF. (2024). Exchange Rate Assessments: CGER Methodology.',
]

add_multiline_text(slide13, Emu(700000), Emu(4300000), Emu(10792000), Emu(2000000),
                   refs, font_size=13, color=GRAY)

# ============================================================
# 第14页：结束页
# ============================================================
slide14 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide14, WHITE)
add_accent_bar(slide14, Emu(0), Emu(0), Emu(12192000), Emu(300000), ACCENT_BLUE)

add_shape_bg(slide14, Emu(0), Emu(2000000), Emu(12192000), Emu(3000000), LIGHT_GRAY)

add_text_box(slide14, Emu(1000000), Emu(2300000), Emu(10192000), Emu(800000),
             '感谢聆听！', font_size=44, color=DARK_BLUE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide14, Emu(1000000), Emu(3200000), Emu(10192000), Emu(600000),
             '购买力平价检验：德美日英实际汇率分析', font_size=24, color=BLUE, alignment=PP_ALIGN.CENTER)

add_text_box(slide14, Emu(1000000), Emu(4200000), Emu(10192000), Emu(400000),
             '📊 数据日期：2026年6月2日 | 📈 6张数据图表 | 📋 完整理论检验',
             font_size=16, color=GRAY, alignment=PP_ALIGN.CENTER)

add_accent_bar(slide14, Emu(0), Emu(6558000), Emu(12192000), Emu(300000), ACCENT_BLUE)

# ============================================================
# 保存
# ============================================================
output_path = '/home/admin/.openclaw/workspace/购买力平价检验分析_演示文稿.pptx'
prs.save(output_path)
print(f'PPT已保存至：{output_path}')
print(f'共 {len(prs.slides)} 张幻灯片')
print('完成！')
