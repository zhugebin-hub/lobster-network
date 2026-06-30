#!/usr/bin/env python3
"""
《商品销售》PPT课件生成器 - 加强版
重点章节：项目1售前准备、项目3推介商品、项目5售后服务
作者：诸葛虾 AI助手
日期：2026-04-25
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os, zipfile

COLORS = {
    'primary': RGBColor(0x1E, 0x88, 0xE5),
    'secondary': RGBColor(0x43, 0xA0, 0x47),
    'accent': RGBColor(0xFF, 0x6F, 0x00),
    'warning': RGBColor(0xE5, 0x39, 0x35),
    'dark': RGBColor(0x26, 0x32, 0x38),
    'light': RGBColor(0xFF, 0xFF, 0xFF),
    'gray': RGBColor(0x78, 0x90, 0x9C),
    'bg': RGBColor(0xF5, 0xF7, 0xFA),
    'purple': RGBColor(0x7B, 0x1F, 0xA2),
    'teal': RGBColor(0x00, 0x96, 0x88),
}

def bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color

def rect(slide, l, t, w, h, c, txt="", fs=14, fc=COLORS['light'], b=False, al=PP_ALIGN.LEFT):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = c; s.line.fill.background()
    if txt:
        tf = s.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.text = txt; p.font.size = Pt(fs); p.font.color.rgb = fc; p.font.bold = b; p.alignment = al
        p.space_before = Pt(4); p.space_after = Pt(4)
    return s

def tb(slide, l, t, w, h, txt, fs=14, fc=COLORS['dark'], b=False, al=PP_ALIGN.LEFT):
    bx = slide.shapes.add_textbox(l, t, w, h)
    tf = bx.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = txt; p.font.size = Pt(fs); p.font.color.rgb = fc; p.font.bold = b; p.alignment = al
    p.font.name = "微软雅黑"
    return bx

def blist(slide, l, t, w, h, items, fs=13, fc=COLORS['dark'], sp=Pt(6)):
    bx = slide.shapes.add_textbox(l, t, w, h)
    tf = bx.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item; p.font.size = Pt(fs); p.font.color.rgb = fc; p.font.name = "微软雅黑"; p.space_after = sp
    return bx

def header(slide, title):
    rect(slide, Inches(0), Inches(0), Inches(10), Inches(1.1), COLORS['primary'])
    tb(slide, Inches(0.5), Inches(0.15), Inches(9), Inches(0.7), title, fs=26, fc=COLORS['light'], b=True)

def proj_cover(slide, num, name, tasks):
    bg(slide, COLORS['primary'])
    rect(slide, Inches(0), Inches(3.2), Inches(10), Inches(0.04), COLORS['light'])
    rect(slide, Inches(0), Inches(5.8), Inches(10), Inches(0.04), COLORS['light'])
    tb(slide, Inches(1), Inches(1.5), Inches(8), Inches(0.8), f"项目{num}", fs=18, fc=RGBColor(0xBB,0xDE,0xFB), b=True)
    tb(slide, Inches(1), Inches(2.2), Inches(8), Inches(1.0), name, fs=36, fc=COLORS['light'], b=True)
    tb(slide, Inches(1), Inches(3.6), Inches(8), Inches(2.0), tasks, fs=14, fc=RGBColor(0xE3,0xF2,0xFD))
    tb(slide, Inches(1), Inches(6.0), Inches(8), Inches(0.5), "《商品销售》· 高等教育出版社 · 中职商贸类专业", fs=11, fc=RGBColor(0x90,0xCA,0xF9))

def content_page(slide, title, items, diff="基础", interactive=False, event=""):
    bg(slide, COLORS['bg']); header(slide, title)
    y = 1.4
    dc = {"基础": COLORS['secondary'], "进阶": COLORS['accent'], "拓展": COLORS['warning']}
    rect(slide, Inches(0.5), Inches(y), Inches(1.3), Inches(0.3), dc.get(diff, COLORS['gray']))
    tb(slide, Inches(0.5), Inches(y), Inches(1.3), Inches(0.3), f"📊 {diff}", fs=9, fc=COLORS['light'], b=True, al=PP_ALIGN.CENTER)
    if interactive:
        rect(slide, Inches(2.0), Inches(y), Inches(1.3), Inches(0.3), COLORS['accent'])
        tb(slide, Inches(2.0), Inches(y), Inches(1.3), Inches(0.3), "💬 互动", fs=9, fc=COLORS['light'], b=True, al=PP_ALIGN.CENTER)
    if event:
        rect(slide, Inches(3.5), Inches(y), Inches(2.5), Inches(0.3), COLORS['warning'])
        tb(slide, Inches(3.5), Inches(y), Inches(2.5), Inches(0.3), "📰 时事热点", fs=9, fc=COLORS['light'], b=True, al=PP_ALIGN.CENTER)
    blist(slide, Inches(0.5), Inches(1.9), Inches(9), Inches(5.0), items, fs=13)

def interact_page(slide, title, q, opts, ans, exp):
    bg(slide, COLORS['bg']); header(slide, title)
    tb(slide, Inches(0.5), Inches(1.3), Inches(9), Inches(0.5), "🤔 思考题", fs=16, fc=COLORS['primary'], b=True)
    tb(slide, Inches(0.5), Inches(1.8), Inches(9), Inches(0.7), q, fs=14, fc=COLORS['dark'])
    y = 2.7
    for i, o in enumerate(opts):
        c = COLORS['primary'] if i % 2 == 0 else COLORS['secondary']
        rect(slide, Inches(0.5), Inches(y + i * 0.45), Inches(9), Inches(0.4), c)
        tb(slide, Inches(0.7), Inches(y + i * 0.45), Inches(8.6), Inches(0.4), o, fs=13, fc=COLORS['light'])
    ay = y + len(opts) * 0.45 + 0.2
    rect(slide, Inches(0.5), Inches(ay), Inches(9), Inches(0.9), COLORS['secondary'])
    tb(slide, Inches(0.7), Inches(ay), Inches(8.6), Inches(0.35), f"✅ 答案：{ans}", fs=13, fc=COLORS['light'], b=True)
    tb(slide, Inches(0.7), Inches(ay + 0.35), Inches(8.6), Inches(0.5), exp, fs=11, fc=RGBColor(0xE8,0xF5,0xE9))

def case_page(slide, title, ct, cc, pts):
    bg(slide, COLORS['bg']); header(slide, title)
    rect(slide, Inches(0.5), Inches(1.3), Inches(9), Inches(0.45), COLORS['accent'])
    tb(slide, Inches(0.7), Inches(1.3), Inches(8.6), Inches(0.45), f"📋 {ct}", fs=15, fc=COLORS['light'], b=True)
    tb(slide, Inches(0.5), Inches(1.9), Inches(9), Inches(1.5), cc, fs=12, fc=COLORS['dark'])
    tb(slide, Inches(0.5), Inches(3.5), Inches(4), Inches(0.3), "💡 分析要点：", fs=13, fc=COLORS['primary'], b=True)
    blist(slide, Inches(0.5), Inches(3.8), Inches(9), Inches(3.0), pts, fs=12)

def summary_page(slide, title, pts, hw=""):
    bg(slide, COLORS['primary'])
    tb(slide, Inches(1), Inches(0.8), Inches(8), Inches(0.7), title, fs=28, fc=COLORS['light'], b=True, al=PP_ALIGN.CENTER)
    rect(slide, Inches(2), Inches(1.6), Inches(6), Inches(0.04), COLORS['light'])
    blist(slide, Inches(1.5), Inches(2.0), Inches(7), Inches(3.5), pts, fs=15, fc=COLORS['light'], sp=Pt(10))
    if hw:
        rect(slide, Inches(1), Inches(5.5), Inches(8), Inches(0.7), RGBColor(0x15,0x65,0xC0))
        tb(slide, Inches(1.2), Inches(5.6), Inches(7.6), Inches(0.5), f"📝 课后作业：{hw}", fs=13, fc=COLORS['light'])

def divider_page(slide, num, name):
    """任务分隔页"""
    bg(slide, COLORS['teal'])
    rect(slide, Inches(0), Inches(2.8), Inches(10), Inches(0.04), COLORS['light'])
    tb(slide, Inches(1), Inches(1.5), Inches(8), Inches(0.8), f"任务{num}", fs=16, fc=RGBColor(0xB2,0xDF,0xDB), b=True)
    tb(slide, Inches(1), Inches(2.2), Inches(8), Inches(1.0), name, fs=32, fc=COLORS['light'], b=True)
    tb(slide, Inches(1), Inches(3.2), Inches(8), Inches(0.5), "《商品销售》· 高等教育出版社", fs=11, fc=RGBColor(0x80,0xCB,0xC4))

# ==================== 项目1：售前准备（加强版）====================

def make_proj1_enhanced():
    """项目1售前准备 - 加强版，14页"""
    prs = Presentation()
    prs.slide_width = Inches(10); prs.slide_height = Inches(7.5)
    
    # 项目封面
    s = prs.slides.add_slide(prs.slide_layouts[6])
    proj_cover(s, "1", "售前准备", "任务1.1 仪容礼仪\n任务1.2 仪表礼仪\n任务1.3 仪态礼仪\n任务1.4 售前心态")
    
    # ===== 任务1.1 仪容礼仪 =====
    s = prs.slides.add_slide(prs.slide_layouts[6])
    divider_page(s, "1.1", "仪容礼仪")
    
    s = prs.slides.add_slide(prs.slide_layouts[6])
    content_page(s, "任务1.1 仪容礼仪", [
        "📌 仪容礼仪的定义",
        "  - 仪容：人的容貌、面部、头发的总体形象",
        "  - 销售人员的仪容是\"无声的名片\"",
        "  - 第一印象7秒法则：客户7秒内形成对你的判断",
        "",
        "📌 面部仪容规范",
        "  - 保持面部清洁：洗脸、护肤、控油",
        "  - 男性：胡须修剪干净或每日剃须",
        "  - 女性：淡妆为宜，避免浓妆艳抹",
        "  - 口腔清洁：口气清新，工作前不吃异味食物",
        "  - 眼部：眼屎清理干净，眼镜保持清洁",
        "",
        "📌 头发规范",
        "  - 清洁整齐，无头屑",
        "  - 男性：前不遮眉、侧不掩耳、后不触领",
        "  - 女性：长发需束起或盘起",
        "  - 发色自然，避免夸张染色",
        "  - 定期修剪（建议2-3周一次）",
    ], diff="基础")
    
    # 仪容自查
    s = prs.slides.add_slide(prs.slide_layouts[6])
    interact_page(s, "任务1.1 仪容自查练习",
        "请对照以下清单，给自己做一个仪容评分（每项1-5分），看看哪些需要改进？",
        ["□ 面部清洁 □ 发型整齐 □ 指甲修剪干净", "□ 口气清新 □ 妆容得体 □ 无异味",
         "□ 耳部清洁 □ 眼镜干净 □ 整体精神饱满", "小组互评：同桌互相检查，提出改进建议"],
        "满分45分，40分以上优秀，30-39分合格，30分以下需要改进",
        "仪容礼仪的核心是\"干净、整洁、精神\"。不需要追求完美外貌，但必须做到干净整洁。第一印象形成只需7秒，仪容是销售成功的第一步。")
    
    # ===== 任务1.2 仪表礼仪 =====
    s = prs.slides.add_slide(prs.slide_layouts[6])
    divider_page(s, "1.2", "仪表礼仪")
    
    s = prs.slides.add_slide(prs.slide_layouts[6])
    content_page(s, "任务1.2 仪表礼仪", [
        "📌 仪表礼仪的定义",
        "  - 仪表：穿着打扮的整体形象",
        "  - \"人靠衣装\"——穿着反映专业度",
        "  - 客户通过穿着判断你的专业水平",
        "",
        "📌 着装基本原则（TPO原则）",
        "  T - Time（时间）：季节、时段",
        "  P - Place（地点）：门店、展会、拜访",
        "  O - Occasion（场合）：日常销售、重要客户接待",
        "",
        "📌 销售人员标准着装",
        "  - 统一制服：整洁、无褶皱、无破损",
        "  - 工牌佩戴：左胸位置，端正清晰",
        "  - 鞋子：深色皮鞋/黑色平底鞋，干净无灰尘",
        "  - 配饰：简约为主，避免夸张首饰",
        "  - 女性：裙长及膝，不穿超短裙",
        "",
        "📌 着装禁忌",
        "  ❌ 过于暴露、过于休闲、过于花哨",
        "  ❌ 运动鞋配正装、拖鞋上岗",
        "  ❌ 衣服有异味、污渍、起球",
    ], diff="基础", interactive=True)
    
    # 案例：胖东来
    s = prs.slides.add_slide(prs.slide_layouts[6])
    case_page(s, "任务1.2 案例分析",
        "案例：胖东来的着装标准",
        "胖东来对员工的着装要求极为严格：制服每日更换、工牌统一佩戴、鞋子必须为黑色。2025年胖东来调改帮扶其他超市时，首先输出的就是\"着装标准化\"方案。其员工形象已成为品牌识别的重要组成部分。",
        ["统一着装提升品牌专业形象",
         "细节决定品质：工牌、鞋子、制服整洁度",
         "着装标准需要制度化、日常检查",
         "2026年零售行业竞争加剧，员工形象成为差异化竞争力"])
    
    # ===== 任务1.3 仪态礼仪 =====
    s = prs.slides.add_slide(prs.slide_layouts[6])
    divider_page(s, "1.3", "仪态礼仪")
    
    s = prs.slides.add_slide(prs.slide_layouts[6])
    content_page(s, "任务1.3 仪态礼仪", [
        "📌 仪态礼仪的定义",
        "  - 仪态：人的姿态、动作、表情",
        "  - \"站有站相、坐有坐相\"——仪态体现修养",
        "  - 仪态比语言更有说服力",
        "",
        "📌 站姿",
        "  - 挺胸收腹，双肩自然放松",
        "  - 双手自然下垂或交叠于腹前",
        "  - 双脚与肩同宽，不倚靠货架/柜台",
        "  - 禁忌：双手抱胸、插兜、抖腿",
        "",
        "📌 微笑服务",
        "  - 微笑是最好的\"破冰工具\"",
        "  - 标准：露出上排8颗牙齿",
        "  - 眼神配合：微笑+眼神交流",
        "  - 练习：咬筷子练微笑（经典方法）",
        "",
        "📌 手势与引导",
        "  - 指引方向：手掌向上，四指并拢",
        "  - 递接物品：双手递接，正面朝向对方",
        "  - 禁忌：用手指指人、单手递物",
    ], diff="基础", interactive=True)
    
    # 仪态模拟
    s = prs.slides.add_slide(prs.slide_layouts[6])
    interact_page(s, "任务1.3 仪态模拟练习",
        "两人一组，一人扮演销售、一人扮演顾客，模拟以下场景并互相评价仪态：",
        ["场景1：顾客进店，销售站立迎接（站姿+微笑+问候）",
         "场景2：顾客询问商品，销售手势引导（手势+眼神）",
         "场景3：递送商品给顾客（双手递接+微笑）",
         "场景4：顾客离开，销售送别（鞠躬+道别）"],
        "评价标准：站姿端正✅ 微笑自然✅ 手势规范✅ 眼神交流✅",
        "仪态礼仪需要\"刻意练习\"。建议每天课前5分钟做仪态训练，形成肌肉记忆。好的仪态不是\"装出来的\"，而是\"练出来的\"。")
    
    # ===== 任务1.4 售前心态 =====
    s = prs.slides.add_slide(prs.slide_layouts[6])
    divider_page(s, "1.4", "售前心态")
    
    s = prs.slides.add_slide(prs.slide_layouts[6])
    content_page(s, "任务1.4 售前心态", [
        "📌 为什么心态比技巧更重要？",
        "  - 技巧可以学，心态决定能不能坚持",
        "  - 销售是\"被拒绝的艺术\"——没有好心态会崩溃",
        "  - 客户能感受到你的情绪：自信vs紧张、热情vs敷衍",
        "",
        "📌 销售人员必备的5种心态",
        "  1️⃣ 自信心态：相信自己的产品，相信自己能做好",
        "  2️⃣ 空杯心态：放下成见，持续学习",
        "  3️⃣ 积极心态：把拒绝当机会，把挫折当成长",
        "  4️⃣ 利他心态：真正为客户着想，不是\"只想卖出去\"",
        "  5️⃣ 韧性心态：被拒绝10次，第11次依然微笑",
        "",
        "📌 心态调整方法",
        "  - 每日晨会：自我激励+团队打气",
        "  - 成功日记：记录每天的小成就",
        "  - 压力释放：运动、倾诉、深呼吸",
        "  - 2026年心理健康被纳入企业员工关怀重点",
    ], diff="进阶", event="2026年人社部将\"销售心理疏导\"纳入职业技能培训目录")
    
    # 心态分享
    s = prs.slides.add_slide(prs.slide_layouts[6])
    interact_page(s, "任务1.4 心态分享",
        "分享一次你被拒绝/失败的经历：当时是什么感受？后来怎么调整的？",
        ["经历1：第一次推销被拒绝，觉得丢脸→后来明白\"拒绝是常态\"",
         "经历2：业绩垫底，想放弃→找到方法后逆袭",
         "经历3：被客户误解投诉→学会\"先处理情绪再处理问题\"",
         "讨论：你平时用什么方法调整心态？"],
        "每个人的经历都是宝贵的成长素材",
        "销售的本质是\"抗挫折能力\"的比拼。据统计，优秀的销售人员平均被拒绝7次后才能成交。关键不是\"不被拒绝\"，而是\"被拒绝后还能微笑面对下一个客户\"。")
    
    # 项目1总结
    s = prs.slides.add_slide(prs.slide_layouts[6])
    summary_page(s, "项目1 总结", [
        "✅ 仪容：干净、整洁、精神——销售的第一张名片",
        "✅ 仪表：TPO原则——穿着体现专业度",
        "✅ 仪态：站姿、微笑、手势——细节决定品质",
        "✅ 心态：自信、积极、韧性——销售的底层能力",
        "✅ 售前准备 = 外在形象 + 内在心态",
    ], hw="录制一段1分钟的视频：展示你的仪容仪表仪态（站姿+微笑+问候+手势引导），自评并写出改进计划")
    
    return prs

# ==================== 项目3：推介商品（加强版）====================

def make_proj3_enhanced():
    """项目3推介商品 - 加强版，14页"""
    prs = Presentation()
    prs.slide_width = Inches(10); prs.slide_height = Inches(7.5)
    
    s = prs.slides.add_slide(prs.slide_layouts[6])
    proj_cover(s, "3", "推介商品", "任务3.1 认识商品推介\n任务3.2 巧用电话推介\n任务3.3 活用新媒体推介")
    
    # ===== 任务3.1 认识商品推介 =====
    s = prs.slides.add_slide(prs.slide_layouts[6])
    divider_page(s, "3.1", "认识商品推介")
    
    s = prs.slides.add_slide(prs.slide_layouts[6])
    content_page(s, "任务3.1 认识商品推介", [
        "📌 什么是商品推介？",
        "  - 将商品的特点、优势、价值传递给客户",
        "  - 不是\"背参数\"，而是\"讲利益\"",
        "  - 推介 = 翻译：把产品语言翻译成客户语言",
        "",
        "📌 FABE推介法（核心方法）",
        "  F - Features（特点）：商品有什么",
        "  A - Advantages（优势）：比别人好在哪",
        "  B - Benefits（利益）：对客户有什么用",
        "  E - Evidence（证据）：怎么证明",
        "",
        "📌 FABE示例：以保温杯为例",
        "  F：\"316不锈钢内胆，保温12小时\"",
        "  A：\"比普通304更耐腐蚀，保温更持久\"",
        "  B：\"早上装的热水，到晚上还是烫的\"",
        "  E：\"这是质检报告，客户好评率98%\"",
        "",
        "📌 推介的\"三多三少\"原则",
        "  多讲利益，少讲参数 | 多问需求，少说产品 | 多给证据，少做承诺",
    ], diff="进阶")
    
    # FABE练习
    s = prs.slides.add_slide(prs.slide_layouts[6])
    interact_page(s, "任务3.1 FABE实战练习",
        "请用FABE法推介以下商品（任选其一，小组内展示）：",
        ["A. 无线蓝牙耳机", "B. 智能手表", "C. 电动牙刷", "D. 自热火锅"],
        "参考框架：F→A→B→E 四句话，每句不超过20字",
        "FABE的精髓是\"把产品语言翻译成客户语言\"。客户不关心\"40dB降噪\"，关心的是\"地铁上也能安静听歌\"。好的推介让客户秒懂\"这对我有什么用\"。")
    
    # FABE进阶练习
    s = prs.slides.add_slide(prs.slide_layouts[6])
    content_page(s, "任务3.1 FABE进阶：多场景推介", [
        "📌 同一商品，不同客户，不同FABE",
        "",
        "📱 商品：智能手表",
        "  对年轻人：F-运动监测 A-精准记录跑步数据 B-帮你科学健身 C-运动APP数据对比",
        "  对中年人：F-心率监测 A-异常及时提醒 B-守护心脏健康 C-三甲医院临床验证",
        "  对老年人：F-跌倒检测 A-自动拨打紧急电话 B-关键时刻救命 C-已挽救200+生命",
        "",
        "📌 推介的\"场景化\"技巧",
        "  - 把参数变成场景：\"防水50米\"→\"游泳洗澡都能戴\"",
        "  - 把功能变成故事：\"续航7天\"→\"出差一周不用带充电器\"",
        "  - 把优势变成对比：\"轻薄\"→\"比名片还薄，放口袋无感\"",
    ], diff="拓展")
    
    # ===== 任务3.2 巧用电话推介 =====
    s = prs.slides.add_slide(prs.slide_layouts[6])
    divider_page(s, "3.2", "巧用电话推介")
    
    s = prs.slides.add_slide(prs.slide_layouts[6])
    content_page(s, "任务3.2 巧用电话推介", [
        "📌 电话推介的特点",
        "  - 看不见表情，全靠声音传递信息",
        "  - 客户防备心强（骚扰电话多）",
        "  - 时间窗口短：前15秒决定去留",
        "",
        "📌 电话推介的标准流程",
        "  1. 开场白：\"您好，我是XX店的XX，打扰您1分钟\"",
        "  2. 表明来意：\"您之前关注过XX产品，现在有活动\"",
        "  3. 利益点：\"比平时便宜80元，还送赠品\"",
        "  4. 引导行动：\"您今天方便来门店看看吗？\"",
        "  5. 结束语：\"好的，期待您的光临，再见\"",
        "",
        "📌 电话推介的技巧",
        "  - 语速适中（每分钟180-200字）",
        "  - 面带微笑（声音能\"听\"到微笑）",
        "  - 站着打电话（声音更有力量）",
        "  - 准备话术脚本，但不要\"念稿\"",
        "",
        "📌 电话推介的禁忌",
        "  ❌ 不报身份直接推销 ❌ 语速过快 ❌ 被客户拒绝后态度变化",
    ], diff="进阶", interactive=True)
    
    # 电话模拟
    s = prs.slides.add_slide(prs.slide_layouts[6])
    interact_page(s, "任务3.2 电话模拟练习",
        "两人一组，用电话模拟以下场景（每人轮流扮演销售和顾客）：",
        ["场景1：给老客户打电话，通知新品到货",
         "场景2：给咨询过的客户打电话，跟进购买意向",
         "场景3：给沉睡客户打电话，推送优惠活动",
         "评价标准：开场白✅ 利益点✅ 引导✅ 应对拒绝✅"],
        "电话推介的成功率通常只有5%-10%，但这很正常",
        "电话推介的核心不是\"一次成功\"，而是\"持续优化\"。建议每次打完电话记录：客户问了什么、拒绝了什么、什么话术有效。积累100通电话后，你会发现自己的话术越来越精准。")
    
    # 电话话术模板
    s = prs.slides.add_slide(prs.slide_layouts[6])
    content_page(s, "任务3.2 电话话术模板库", [
        "📌 模板1：老客户回访型",
        "  \"王姐您好，我是XX服装店的小李。您上次买的连衣裙穿得还满意吗？...太好了！告诉您个好消息，我们刚到了一批新款，特别适合您，今天有空来试试吗？\"",
        "",
        "📌 模板2：活动通知型",
        "  \"您好，我是XX电器的小张。您之前咨询过的XX冰箱，这周末有特价活动，比平时便宜300元，还送延保一年。名额有限，我帮您留一个？\"",
        "",
        "📌 模板3：沉睡客户激活型",
        "  \"您好，我是XX美妆的小陈。好久没看到您了，最近我们推出了会员专属优惠，满200减50，还有新品体验装免费送。您这周方便过来吗？\"",
        "",
        "📌 应对拒绝的话术",
        "  \"没关系，感谢您的时间。这是我的微信，有需要随时联系我。\"",
    ], diff="进阶")
    
    # ===== 任务3.3 活用新媒体推介 =====
    s = prs.slides.add_slide(prs.slide_layouts[6])
    divider_page(s, "3.3", "活用新媒体推介")
    
    s = prs.slides.add_slide(prs.slide_layouts[6])
    content_page(s, "任务3.3 活用新媒体推介", [
        "📌 新媒体推介的渠道",
        "  - 微信朋友圈/视频号：日常种草",
        "  - 抖音/快手短视频：产品展示+使用场景",
        "  - 直播带货：实时互动+限时优惠",
        "  - 小红书：种草笔记+真实体验",
        "  - 企业微信社群：精准推送+社群运营",
        "",
        "📌 短视频推介要点",
        "  - 前3秒抓眼球：痛点/反差/悬念",
        "  - 15-30秒讲清卖点：一个视频一个卖点",
        "  - 结尾引导行动：\"点击购物车\"\"私信咨询\"",
        "  - 2026年短视频带货占零售销售额的35%",
        "",
        "📌 直播推介要点",
        "  - 人设：专业、真实、有亲和力",
        "  - 节奏：5分钟一个产品，穿插互动",
        "  - 话术：\"宝宝们\"\"最后10单\"\"限时秒杀\"",
        "  - 2026年AI数字人直播成为新趋势",
    ], diff="拓展", event="2026年商务部发布《数字商务高质量发展行动计划》，新媒体销售纳入职业技能标准")
    
    # 新媒体策划
    s = prs.slides.add_slide(prs.slide_layouts[6])
    interact_page(s, "任务3.3 新媒体推介策划",
        "选择一款商品，设计一个新媒体推介方案（小组合作）：",
        ["方案A：拍摄一条15秒短视频（写脚本+分镜头）",
         "方案B：策划一场30分钟直播（写流程+话术）",
         "方案C：写一篇小红书种草笔记（标题+正文+配图）",
         "方案D：设计一个朋友圈营销文案（3条系列）"],
        "评价标准：吸引力✅ 信息量✅ 行动引导✅ 合规性✅",
        "新媒体推介的核心是\"内容即广告\"。好的内容让客户主动想看、主动分享。2026年AI工具可以辅助生成文案、剪辑视频，但\"网感\"和\"创意\"仍然需要人工。")
    
    # 新媒体案例
    s = prs.slides.add_slide(prs.slide_layouts[6])
    case_page(s, "任务3.3 案例分析",
        "案例：东方甄选的知识型直播",
        "东方甄选以\"知识+带货\"模式脱颖而出。主播在推介商品时，不仅讲产品，还讲文化、讲历史、讲英语。2025年其GMV突破100亿。其核心逻辑是：内容吸引用户→信任促进购买→口碑带来复购。",
        ["新媒体推介的核心是\"内容价值\"而非\"硬广\"",
         "知识型内容建立信任，信任促进转化",
         "2026年AI可以辅助生成内容，但\"人设\"和\"温度\"不可替代",
         "新媒体推介需要\"网感\"——理解平台调性和用户喜好"])
    
    # 项目3总结
    s = prs.slides.add_slide(prs.slide_layouts[6])
    summary_page(s, "项目3 总结", [
        "✅ FABE推介：把产品特点翻译成客户利益",
        "✅ 电话推介：15秒黄金开场，声音传递信任",
        "✅ 新媒体推介：短视频+直播+社群=新三件套",
        "✅ 推介的核心：不是\"卖出去\"，而是\"帮到客户\"",
        "✅ 时代趋势：AI辅助+人工创意=最佳组合",
    ], hw="选择一款商品，制作一份新媒体推介方案（短视频脚本/直播策划/朋友圈文案，任选其一）")
    
    return prs

# ==================== 项目5：售后服务（加强版）====================

def make_proj5_enhanced():
    """项目5售后服务 - 加强版，14页"""
    prs = Presentation()
    prs.slide_width = Inches(10); prs.slide_height = Inches(7.5)
    
    s = prs.slides.add_slide(prs.slide_layouts[6])
    proj_cover(s, "5", "售后服务", "任务5.1 售后接待\n任务5.2 \"三包\"服务\n任务5.3 投诉处理\n任务5.4 信息整理与反馈")
    
    # ===== 任务5.1 售后接待 =====
    s = prs.slides.add_slide(prs.slide_layouts[6])
    divider_page(s, "5.1", "售后接待")
    
    s = prs.slides.add_slide(prs.slide_layouts[6])
    content_page(s, "任务5.1 售后接待", [
        "📌 售后接待的重要性",
        "  - 售后是客户体验的\"最后一公里\"",
        "  - 好售后 = 复购 + 转介绍",
        "  - 维护老客户成本 = 开发新客户的1/5",
        "  - 客户满意度提升5%，利润增加25%-85%（贝恩咨询）",
        "",
        "📌 售后接待的标准流程",
        "  1. 热情接待：\"您好，有什么可以帮您？\"",
        "  2. 了解问题：\"请具体说一下情况...\"",
        "  3. 确认信息：查看凭证、检查商品",
        "  4. 提出方案：\"我给您...方案，您看可以吗？\"",
        "  5. 执行处理：退换货/维修/补偿",
        "  6. 回访确认：\"问题解决了吗？还有其他需要吗？\"",
        "",
        "📌 售后接待的态度",
        "  - 不推诿、不敷衍、不冷漠",
        "  - 即使不是我们的责任，也要\"先解决情绪，再解决问题\"",
        "  - 把售后接待当作\"二次销售\"的机会",
    ], diff="基础")
    
    # 售后接待模拟
    s = prs.slides.add_slide(prs.slide_layouts[6])
    interact_page(s, "任务5.1 售后接待模拟",
        "场景：一位客户拿着一个月前买的衣服来退货，说\"穿了两次就不想穿了\"。但按规定超过7天不支持无理由退货。请模拟售后接待过程。",
        ["第一步：热情接待，了解情况",
         "第二步：查看凭证，确认是否符合政策",
         "第三步：解释政策，同时提供替代方案",
         "第四步：给出解决方案（如：换货/积分补偿/优惠券）"],
        "参考答案：先共情→再解释政策→最后给替代方案",
        "售后接待的核心原则：\"制度要有底线，服务要有温度\"。即使不能满足客户的要求，也要让客户感受到被尊重。可以拒绝\"退货\"，但不能拒绝\"帮助\"。")
    
    # ===== 任务5.2 \"三包\"服务 =====
    s = prs.slides.add_slide(prs.slide_layouts[6])
    divider_page(s, "5.2", "\"三包\"服务")
    
    s = prs.slides.add_slide(prs.slide_layouts[6])
    content_page(s, "任务5.2 \"三包\"服务", [
        "📌 什么是\"三包\"？",
        "  - 包修、包换、包退——国家法定要求",
        "  - 保护消费者权益，规范市场秩序",
        "",
        "📌 \"三包\"基本规定",
        "  🔧 包修：产品出现问题，免费修理",
        "  🔄 包换：符合换货条件的，免费更换同型号",
        "  ⏪ 包退：符合退货条件的，全额退款",
        "",
        "📌 不同品类的\"三包\"期限",
        "  📱 电子产品：7天退货，15天换货，1年保修",
        "  👕 服装鞋帽：7天质量问题退换",
        "  🏠 家用电器：7天退货，15天换货，整机1年/主要部件3年",
        "  🚗 汽车：7天退货，15天换货，2年/5万公里保修",
        "",
        "📌 \"三包\"注意事项",
        "  - 保留购买凭证（发票、小票）",
        "  - 人为损坏不在\"三包\"范围内",
        "  - 超过\"三包\"期限可收费维修",
        "  - 2026年《消费者权益保护法》修订，部分品类\"三包\"期限延长",
    ], diff="基础")
    
    # 三包案例
    s = prs.slides.add_slide(prs.slide_layouts[6])
    case_page(s, "任务5.2 案例分析",
        "案例：手机\"三包\"纠纷",
        "消费者小王购买手机10天后发现屏幕闪烁，到门店要求退货。门店表示超过7天只能维修不能退货。小王投诉到12315，经调解，门店同意换货（非退货），因为屏幕闪烁属于质量问题，符合15天换货条件。",
        ["\"三包\"政策因品类而异，销售人员必须熟悉",
         "质量问题vs人为损坏：界定标准要清楚",
         "遇到\"三包\"争议，及时上报，不要自行决定",
         "2026年消费者权益保护力度加大，企业需更重视\"三包\"合规"])
    
    # ===== 任务5.3 投诉处理 =====
    s = prs.slides.add_slide(prs.slide_layouts[6])
    divider_page(s, "5.3", "投诉处理")
    
    s = prs.slides.add_slide(prs.slide_layouts[6])
    content_page(s, "任务5.3 投诉处理", [
        "📌 投诉处理的心态",
        "  - 投诉是礼物——客户给你改进的机会",
        "  - 投诉处理得好 = 客户忠诚度提升",
        "  - 一个满意的投诉处理 = 免费广告",
        "",
        "📌 HEARD投诉处理法（迪士尼标准）",
        "  H - Hear 倾听：让客户说完，不打断",
        "  E - Empathize 共情：\"我理解您的感受\"",
        "  A - Apologize 道歉：为不好的体验道歉",
        "  R - Resolve 解决：提供方案，给客户选择",
        "  D - Diagnose 诊断：事后分析，预防再发",
        "",
        "📌 投诉处理中的沟通技巧",
        "  - 控制情绪，不被客户的愤怒影响",
        "  - 用\"我\"代替\"你\"（\"我来帮您\"vs\"你应该\"）",
        "  - 给出明确时间承诺（\"24小时内回复\"）",
        "  - 不轻易说\"不\"，说\"让我帮您想办法\"",
    ], diff="进阶", interactive=True, event="2026年央视315晚会：\"公平消费 诚信守护\"，投诉处理成为焦点")
    
    # 投诉处理模拟
    s = prs.slides.add_slide(prs.slide_layouts[6])
    interact_page(s, "任务5.3 投诉处理实战",
        "客户怒气冲冲：\"你们这个洗衣机买了才一个月就坏了，耽误了我多少事！我要投诉！\"请用HEARD法处理。",
        ["H：耐心听完，记录关键信息（不打断）",
         "E：\"非常理解您的心情，洗衣机坏了确实很耽误事\"",
         "A：\"给您带来不便，我代表门店向您道歉\"",
         "R：\"我给您两个方案：①立即换新 ②先借一台给您用，这台马上维修\""],
        "参考答案：完整走一遍HEARD流程",
        "处理投诉的核心：先处理情绪，再处理问题。客户愤怒时讲道理是无效的。先共情道歉，让客户情绪降下来，再给解决方案。给客户选择权，让客户感觉\"被尊重\"。")
    
    # 投诉升级处理
    s = prs.slides.add_slide(prs.slide_layouts[6])
    content_page(s, "任务5.3 投诉升级处理", [
        "📌 投诉分级处理机制",
        "  🟢 一级投诉（一线员工）：",
        "    - 小额退换货、一般咨询",
        "    - 权限：50元以内补偿",
        "",
        "  🟡 二级投诉（店长/主管）：",
        "    - 质量问题争议、服务态度投诉",
        "    - 权限：200元以内补偿/换货",
        "",
        "  🔴 三级投诉（公司层面）：",
        "    - 重大质量事故、媒体曝光、12315投诉",
        "    - 权限：全额退款/召回/法务介入",
        "",
        "📌 投诉升级的原则",
        "  - 一线能解决的，不上交",
        "  - 超出权限的，及时上报",
        "  - 重大投诉，24小时内响应",
        "  - 所有投诉，必须闭环（有始有终）",
    ], diff="拓展")
    
    # ===== 任务5.4 信息整理与反馈 =====
    s = prs.slides.add_slide(prs.slide_layouts[6])
    divider_page(s, "5.4", "信息整理与反馈")
    
    s = prs.slides.add_slide(prs.slide_layouts[6])
    content_page(s, "任务5.4 信息整理与反馈", [
        "📌 为什么要整理售后信息？",
        "  - 售后数据是产品改进的\"金矿\"",
        "  - 客户反馈能发现产品设计的盲点",
        "  - 数据驱动决策，不是\"凭感觉\"",
        "",
        "📌 售后信息整理的内容",
        "  📊 退换货数据：数量、原因、品类分布",
        "  📊 投诉数据：类型、频率、处理结果",
        "  📊 客户建议：产品改进、服务优化",
        "  📊 满意度数据：评分、评价、NPS",
        "",
        "📌 反馈机制",
        "  - 日记录：每天记录售后情况",
        "  - 周汇总：每周整理分析，提交店长",
        "  - 月报告：每月形成分析报告，反馈给供应商",
        "  - 紧急反馈：重大质量问题，24小时内上报",
        "",
        "📌 数字化工具",
        "  - CRM系统：自动记录、分类、分析",
        "  - AI客服：自动回复常见问题，人工处理复杂投诉",
        "  - 数据看板：实时展示售后关键指标",
    ], diff="拓展")
    
    # 项目5总结
    s = prs.slides.add_slide(prs.slide_layouts[6])
    summary_page(s, "项目5 总结", [
        "✅ 售后接待：热情、专业、有温度",
        "✅ \"三包\"服务：法定要求，必须熟悉",
        "✅ 投诉处理：HEARD法，先情绪后问题",
        "✅ 信息反馈：数据驱动，持续改进",
        "✅ 售后不是成本，是品牌最值的投资",
        "✅ 好售后 = 复购 + 转介绍 + 口碑",
    ], hw="选择一家你熟悉的企业，分析其售后服务体系（至少包含3个环节），写一份500字的分析报告")
    
    return prs

# ==================== 主程序 ====================

def main():
    output_dir = "/home/admin/.openclaw/workspace/商品销售课件_加强版"
    os.makedirs(output_dir, exist_ok=True)
    
    projects = [
        ("项目1_售前准备_加强版.pptx", make_proj1_enhanced),
        ("项目3_推介商品_加强版.pptx", make_proj3_enhanced),
        ("项目5_售后服务_加强版.pptx", make_proj5_enhanced),
    ]
    
    files = []
    for fname, func in projects:
        print(f"📚 正在生成 {fname}...")
        prs = func()
        fpath = os.path.join(output_dir, fname)
        prs.save(fpath)
        files.append(fpath)
        print(f"  ✅ {fname} — {len(prs.slides)} 页")
    
    # 打包
    zip_file = os.path.join(output_dir, "商品销售课件_三核心章节加强版.zip")
    with zipfile.ZipFile(zip_file, 'w', zipfile.ZIP_DEFLATED) as zf:
        for f in files:
            zf.write(f, os.path.basename(f))
    
    total_slides = 0
    for _, func in projects:
        prs = func()
        total_slides += len(prs.slides)
    
    print(f"\n📦 打包完成：{zip_file}")
    print(f"\n📊 课件统计：")
    print(f"  项目1 售前准备（加强版）：{len(make_proj1_enhanced().slides)} 页")
    print(f"  项目3 推介商品（加强版）：{len(make_proj3_enhanced().slides)} 页")
    print(f"  项目5 售后服务（加强版）：{len(make_proj5_enhanced().slides)} 页")
    print(f"  总计：{total_slides} 页")
    print(f"\n🎯 加强版特色：")
    print(f"  ✅ 每个任务独立分隔页，结构清晰")
    print(f"  ✅ 每任务至少2页内容，讲解更充分")
    print(f"  ✅ 更多互动环节：角色扮演、情景模拟、小组讨论")
    print(f"  ✅ 更多案例分析：胖东来、东方甄选等真实案例")
    print(f"  ✅ 话术模板库：电话推介话术、投诉处理话术")
    print(f"  ✅ 难度递进：基础→进阶→拓展")

if __name__ == "__main__":
    main()
