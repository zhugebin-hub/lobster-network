#!/usr/bin/env python3
"""生成法治演讲稿 PPT"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# === 配色方案 ===
PRIMARY = RGBColor(0x8B, 0x00, 0x00)      # 深红（佛教/法治主题色）
ACCENT = RGBColor(0xD4, 0xA0, 0x17)       # 金色
DARK = RGBColor(0x2C, 0x2C, 0x2C)         # 深灰
LIGHT_BG = RGBColor(0xFA, 0xF8, 0xF5)     # 米白背景
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x66, 0x66, 0x66)

# === 辅助函数 ===
def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape_bg(slide, color, left=0, top=0, width=None, height=None):
    """添加矩形背景"""
    w = width or prs.slide_width
    h = height or prs.slide_height
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_textbox(slide, left, top, width, height, text, font_size=18, 
                color=DARK, bold=False, alignment=PP_ALIGN.LEFT, font_name='微软雅黑'):
    """添加文本框"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_multi_para(slide, left, top, width, height, lines, font_size=16, 
                   color=DARK, line_spacing=1.5, bold_first=False):
    """添加多行文本"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = '微软雅黑'
        p.space_after = Pt(font_size * line_spacing * 0.4)
        if bold_first and i == 0:
            p.font.bold = True
    return txBox

def add_accent_line(slide, left, top, width, color=PRIMARY, height=Pt(3)):
    """添加装饰线"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_number_circle(slide, left, top, size, number, color=PRIMARY):
    """添加圆形数字标记"""
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = str(number)
    p.font.size = Pt(int(size / Inches(1) * 14))
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = '微软雅黑'
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(0)
    tf.paragraphs[0].space_after = Pt(0)
    return shape

# ============================================================
# 第1页：封面
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白版式
set_slide_bg(slide, RGBColor(0x1A, 0x0A, 0x0A))

# 顶部装饰条
add_shape_bg(slide, PRIMARY, Inches(0), Inches(0), prs.slide_width, Pt(6))

# 主标题
add_textbox(slide, Inches(1.5), Inches(2.0), Inches(10.3), Inches(1.5),
            '以戒为师  依法护心', font_size=44, color=WHITE, bold=True, 
            alignment=PP_ALIGN.CENTER)

# 装饰线
add_accent_line(slide, Inches(4.5), Inches(3.6), Inches(4.3), ACCENT, Pt(2))

# 副标题
add_textbox(slide, Inches(1.5), Inches(3.9), Inches(10.3), Inches(1.0),
            '从佛教戒律精神阐释"法治"核心价值观', font_size=22, color=ACCENT,
            alignment=PP_ALIGN.CENTER)

# 底部信息
add_textbox(slide, Inches(1.5), Inches(5.5), Inches(10.3), Inches(0.8),
            '双通班专题分享', font_size=16, color=GRAY, alignment=PP_ALIGN.CENTER)

# ============================================================
# 第2页：引言 / 开场
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, LIGHT_BG)

# 左侧色块
add_shape_bg(slide, PRIMARY, Inches(0), Inches(0), Inches(0.15), prs.slide_height)

# 标题
add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.6),
            '引言：为什么要谈"法治"？', font_size=32, color=PRIMARY, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(2.5))

# 内容
lines = [
    '在座各位来自五大宗教，各有信仰传统和戒律规范',
    '有人问：我是学佛的，讲戒律就够了，为什么要讲法治？',
    '道教讲"道法自然"，伊斯兰教讲"沙里亚"，天主教讲自然法',
    '每个宗教都有行为规范、遵守法律的教导',
    '',
    '宗教的戒律、教规，与国家的法律之间，到底是什么关系？',
]
add_multi_para(slide, Inches(1.2), Inches(1.8), Inches(10.5), Inches(4.5),
               lines, font_size=20, color=DARK, line_spacing=1.8)

# 底部提示
add_textbox(slide, Inches(1.2), Inches(6.2), Inches(10), Inches(0.5),
            '今天，从佛教视角出发，尝试给出答案', font_size=16, color=ACCENT, bold=True)

# ============================================================
# 第3页：目录
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, LIGHT_BG)

add_shape_bg(slide, PRIMARY, Inches(0), Inches(0), Inches(0.15), prs.slide_height)

add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.6),
            '内容纲要', font_size=32, color=PRIMARY, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(2.5))

toc_items = [
    ('一', '制度传承，高于个人'),
    ('二', '规则面前，人人平等'),
    ('三', '慈悲生起，自觉持戒'),
    ('四', '德法相济，闭环成序'),
    ('五', '尊法守戒，从我做起'),
]

for i, (num, title) in enumerate(toc_items):
    y = Inches(1.8) + Inches(i * 0.95)
    add_number_circle(slide, Inches(1.2), y, Inches(0.55), num, PRIMARY)
    add_textbox(slide, Inches(2.0), y + Inches(0.05), Inches(9), Inches(0.5),
                title, font_size=22, color=DARK)

# ============================================================
# 第4页：第一部分标题页
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, RGBColor(0x1A, 0x0A, 0x0A))

add_textbox(slide, Inches(1.5), Inches(1.5), Inches(10), Inches(0.8),
            '第一部分', font_size=24, color=ACCENT, alignment=PP_ALIGN.CENTER)

add_accent_line(slide, Inches(4.5), Inches(2.5), Inches(4.3), ACCENT, Pt(2))

add_textbox(slide, Inches(1.5), Inches(2.8), Inches(10), Inches(1.2),
            '制度传承，高于个人', font_size=42, color=WHITE, bold=True,
            alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(2), Inches(4.5), Inches(9.3), Inches(1.5),
            '"佛灭度后，以谁为师？"\n——以戒为师', font_size=22, color=ACCENT,
            alignment=PP_ALIGN.CENTER)

# ============================================================
# 第5页：以戒为师
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, LIGHT_BG)

add_shape_bg(slide, PRIMARY, Inches(0), Inches(0), Inches(0.15), prs.slide_height)

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
            '一、制度传承，高于个人', font_size=30, color=PRIMARY, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.0), Inches(2.5))

# 引用
quote_box = add_shape_bg(slide, RGBColor(0xF5, 0xF0, 0xE8), Inches(1.0), Inches(1.5), Inches(11.3), Inches(1.8))
add_textbox(slide, Inches(1.5), Inches(1.7), Inches(10.3), Inches(1.5),
            '"汝等比丘，于我灭后，当尊重珍敬波罗提木叉，\n如暗遇明、贫人得宝。当知此则是汝等大师，\n若我住世无异此也。"', font_size=18, color=DARK,
            alignment=PP_ALIGN.CENTER)

# 核心要点
add_textbox(slide, Inches(1.0), Inches(3.7), Inches(11), Inches(0.5),
            '佛陀的选择：制度传承，而非个人传承', font_size=22, color=PRIMARY, bold=True)

points = [
    '制度的权威高于个人的权威——谁都不能凌驾于规则之上',
    '规则的作用是"如暗遇明"——什么该做、什么不该做，一目了然',
    '制度的有效性靠共同信守——戒律在，就等于佛陀在',
]
add_multi_para(slide, Inches(1.2), Inches(4.3), Inches(10.5), Inches(2.5),
               points, font_size=18, color=DARK, line_spacing=2.0)

# 金句
add_textbox(slide, Inches(1.0), Inches(6.3), Inches(11), Inches(0.6),
            '戒律就是佛门的"法律"，法律就是世间的戒律', font_size=18, color=ACCENT, bold=True,
            alignment=PP_ALIGN.CENTER)

# ============================================================
# 第6页：第二部分标题页
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, RGBColor(0x1A, 0x0A, 0x0A))

add_textbox(slide, Inches(1.5), Inches(1.5), Inches(10), Inches(0.8),
            '第二部分', font_size=24, color=ACCENT, alignment=PP_ALIGN.CENTER)

add_accent_line(slide, Inches(4.5), Inches(2.5), Inches(4.3), ACCENT, Pt(2))

add_textbox(slide, Inches(1.5), Inches(2.8), Inches(10), Inches(1.2),
            '规则面前，人人平等', font_size=42, color=WHITE, bold=True,
            alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(2), Inches(4.5), Inches(9.3), Inches(1.5),
            '依法不依人', font_size=28, color=ACCENT,
            alignment=PP_ALIGN.CENTER)

# ============================================================
# 第7页：依法不依人
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, LIGHT_BG)

add_shape_bg(slide, PRIMARY, Inches(0), Inches(0), Inches(0.15), prs.slide_height)

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
            '二、规则面前，人人平等', font_size=30, color=PRIMARY, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.0), Inches(2.5))

# 四依四不依
add_textbox(slide, Inches(1.0), Inches(1.5), Inches(11), Inches(0.5),
            '"四依四不依" —— 大乘佛教重要原则', font_size=20, color=DARK, bold=True)

add_textbox(slide, Inches(1.5), Inches(2.1), Inches(10), Inches(0.5),
            '第一条：依法不依人 ——《大宝积经》', font_size=18, color=PRIMARY)

# 随方毗尼
add_shape_bg(slide, RGBColor(0xF5, 0xF0, 0xE8), Inches(1.0), Inches(2.8), Inches(11.3), Inches(2.0))
add_textbox(slide, Inches(1.3), Inches(2.9), Inches(10.7), Inches(0.4),
            '随方毗尼 —— 因地制宜的法治智慧', font_size=18, color=PRIMARY, bold=True)
add_textbox(slide, Inches(1.5), Inches(3.4), Inches(10.3), Inches(1.3),
            '"虽是我所制，而于余方不以为清净者，皆不应用；\n虽非我所制，而于余方必应行者，皆不得不行。"\n——《五分律》', font_size=16, color=DARK,
            alignment=PP_ALIGN.CENTER)

# 核心
add_textbox(slide, Inches(1.0), Inches(5.2), Inches(11), Inches(0.5),
            '佛陀留出了因地制宜的空间，但底线从未动摇：', font_size=18, color=DARK)
add_textbox(slide, Inches(1.0), Inches(5.7), Inches(11), Inches(0.5),
            '规则面前人人平等，不可动摇', font_size=22, color=PRIMARY, bold=True,
            alignment=PP_ALIGN.CENTER)

# ============================================================
# 第8页：第三部分标题页
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, RGBColor(0x1A, 0x0A, 0x0A))

add_textbox(slide, Inches(1.5), Inches(1.5), Inches(10), Inches(0.8),
            '第三部分', font_size=24, color=ACCENT, alignment=PP_ALIGN.CENTER)

add_accent_line(slide, Inches(4.5), Inches(2.5), Inches(4.3), ACCENT, Pt(2))

add_textbox(slide, Inches(1.5), Inches(2.8), Inches(10), Inches(1.2),
            '慈悲生起，自觉持戒', font_size=42, color=WHITE, bold=True,
            alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(2), Inches(4.5), Inches(9.3), Inches(1.5),
            '从"不敢"到"不忍"', font_size=28, color=ACCENT,
            alignment=PP_ALIGN.CENTER)

# ============================================================
# 第9页：不敢 vs 不忍
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, LIGHT_BG)

add_shape_bg(slide, PRIMARY, Inches(0), Inches(0), Inches(0.15), prs.slide_height)

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
            '三、慈悲生起，自觉持戒', font_size=30, color=PRIMARY, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.0), Inches(2.5))

# 左右对比
# 左侧：不敢
add_shape_bg(slide, RGBColor(0xF0, 0xF0, 0xF0), Inches(0.8), Inches(1.5), Inches(5.5), Inches(2.5))
add_textbox(slide, Inches(1.0), Inches(1.6), Inches(5.1), Inches(0.5),
            '不敢做', font_size=28, color=GRAY, bold=True, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1.0), Inches(2.2), Inches(5.1), Inches(1.5),
            '因为外面有警察、有监控、有惩罚\n外在约束，怕犯法', font_size=16, color=DARK,
            alignment=PP_ALIGN.CENTER)

# 右侧：不忍
add_shape_bg(slide, RGBColor(0xFD, 0xF5, 0xE6), Inches(7.0), Inches(1.5), Inches(5.5), Inches(2.5))
add_textbox(slide, Inches(7.2), Inches(1.6), Inches(5.1), Inches(0.5),
            '不忍做', font_size=28, color=PRIMARY, bold=True, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(7.2), Inches(2.2), Inches(5.1), Inches(1.5),
            '因为你的心告诉你：那样做会伤害到别人\n内在觉醒，慈悲为怀', font_size=16, color=DARK,
            alignment=PP_ALIGN.CENTER)

# 五戒对照
add_textbox(slide, Inches(0.8), Inches(4.3), Inches(11), Inches(0.5),
            '五戒与法治价值对照', font_size=20, color=PRIMARY, bold=True)

# 表格
table = slide.shapes.add_table(6, 4, Inches(0.8), Inches(4.9), Inches(11.7), Inches(2.3))
table_obj = table.table

headers = ['五戒', '不是怕…', '而是不忍…', '法治价值']
data = [
    ['不杀生', '怕犯戒', '伤害生命', '敬畏生命权'],
    ['不偷盗', '怕被抓', '侵犯他人成果', '尊重财产权'],
    ['不邪淫', '怕名声损', '破坏家庭幸福', '守护婚姻制度'],
    ['不妄语', '怕失信', '欺骗信任者', '维护诚信体系'],
    ['不饮酒', '怕失态', '失去理性伤人', '持守自律精神'],
]

for i, h in enumerate(headers):
    cell = table_obj.cell(0, i)
    cell.text = h
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.font.name = '微软雅黑'
        p.alignment = PP_ALIGN.CENTER
    cell.fill.solid()
    cell.fill.fore_color.rgb = PRIMARY

for r, row_data in enumerate(data):
    for c, val in enumerate(row_data):
        cell = table_obj.cell(r+1, c)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(12)
            p.font.color.rgb = DARK
            p.font.name = '微软雅黑'
            p.alignment = PP_ALIGN.CENTER
        if r % 2 == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0xF8, 0xF6, 0xF3)
        else:
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE

# ============================================================
# 第10页：第四部分标题页
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, RGBColor(0x1A, 0x0A, 0x0A))

add_textbox(slide, Inches(1.5), Inches(1.5), Inches(10), Inches(0.8),
            '第四部分', font_size=24, color=ACCENT, alignment=PP_ALIGN.CENTER)

add_accent_line(slide, Inches(4.5), Inches(2.5), Inches(4.3), ACCENT, Pt(2))

add_textbox(slide, Inches(1.5), Inches(2.8), Inches(10), Inches(1.2),
            '德法相济，闭环成序', font_size=42, color=WHITE, bold=True,
            alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(2), Inches(4.5), Inches(9.3), Inches(1.5),
            '法治为根本 · 德治为辅助', font_size=24, color=ACCENT,
            alignment=PP_ALIGN.CENTER)

# ============================================================
# 第11页：德法相济
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, LIGHT_BG)

add_shape_bg(slide, PRIMARY, Inches(0), Inches(0), Inches(0.15), prs.slide_height)

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
            '四、德法相济，闭环成序', font_size=30, color=PRIMARY, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.0), Inches(2.5))

# 论语引用
add_shape_bg(slide, RGBColor(0xF5, 0xF0, 0xE8), Inches(1.0), Inches(1.5), Inches(11.3), Inches(1.5))
add_textbox(slide, Inches(1.5), Inches(1.6), Inches(10.3), Inches(0.4),
            '《论语》：道之以政，齐之以刑，民免而无耻', font_size=16, color=DARK,
            alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1.5), Inches(2.0), Inches(10.3), Inches(0.8),
            '道之以德，齐之以礼，有耻且格', font_size=16, color=PRIMARY, bold=True,
            alignment=PP_ALIGN.CENTER)

# 核心逻辑
add_textbox(slide, Inches(1.0), Inches(3.4), Inches(11), Inches(0.5),
            '单靠法律 → 只是不敢犯；加上道德教化 → 不忍犯', font_size=20, color=DARK, bold=True,
            alignment=PP_ALIGN.CENTER)

# 三个要点
items = [
    ('法律是底线，是根本', '任何宗教戒律若与国家法律抵触，必须以法律为准'),
    ('宗教是辅助，是升华', '戒律教规帮助信众从内心认同法律的正当性'),
    ('德法相济，闭环成序', '法治为根本，德治为辅助，共同维护社会秩序'),
]

for i, (title, desc) in enumerate(items):
    y = Inches(4.1) + Inches(i * 0.9)
    add_number_circle(slide, Inches(1.2), y, Inches(0.5), i+1, PRIMARY)
    add_textbox(slide, Inches(2.0), y, Inches(9.5), Inches(0.4),
                title, font_size=18, color=DARK, bold=True)
    add_textbox(slide, Inches(2.0), y + Inches(0.4), Inches(9.5), Inches(0.4),
                desc, font_size=14, color=GRAY)

# ============================================================
# 第12页：第五部分标题页
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, RGBColor(0x1A, 0x0A, 0x0A))

add_textbox(slide, Inches(1.5), Inches(1.5), Inches(10), Inches(0.8),
            '第五部分', font_size=24, color=ACCENT, alignment=PP_ALIGN.CENTER)

add_accent_line(slide, Inches(4.5), Inches(2.5), Inches(4.3), ACCENT, Pt(2))

add_textbox(slide, Inches(1.5), Inches(2.8), Inches(10), Inches(1.2),
            '尊法守戒，从我做起', font_size=42, color=WHITE, bold=True,
            alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(2), Inches(4.5), Inches(9.3), Inches(1.5),
            '参与者和建设者', font_size=28, color=ACCENT,
            alignment=PP_ALIGN.CENTER)

# ============================================================
# 第13页：三点践行
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, LIGHT_BG)

add_shape_bg(slide, PRIMARY, Inches(0), Inches(0), Inches(0.15), prs.slide_height)

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
            '五、尊法守戒，从我做起', font_size=30, color=PRIMARY, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.0), Inches(2.5))

add_textbox(slide, Inches(1.0), Inches(1.3), Inches(11), Inches(0.5),
            '作为双通班学员，三点践行：', font_size=18, color=DARK)

# 三个卡片
card_data = [
    ('个人修学', '学法懂法\n做知法守法的表率', '守法本身就是信仰实践的一部分'),
    ('信众引导', '以戒律传统守护法治底线\n善用教义讲清守法即信仰', '每个宗教都有遵守法律的经典依据'),
    ('社会参与', '运用法律维护正当权益\n积极投身公益慈善', '为"中国式现代化"贡献宗教界力量'),
]

for i, (title, content, note) in enumerate(card_data):
    x = Inches(0.8) + Inches(i * 4.1)
    # 卡片背景
    add_shape_bg(slide, WHITE, x, Inches(1.9), Inches(3.8), Inches(4.8))
    # 标题
    add_textbox(slide, x + Inches(0.2), Inches(2.1), Inches(3.4), Inches(0.5),
                title, font_size=20, color=PRIMARY, bold=True, alignment=PP_ALIGN.CENTER)
    add_accent_line(slide, x + Inches(0.5), Inches(2.65), Inches(2.8), ACCENT, Pt(2))
    # 内容
    add_textbox(slide, x + Inches(0.3), Inches(2.9), Inches(3.2), Inches(2.0),
                content, font_size=15, color=DARK, alignment=PP_ALIGN.CENTER)
    # 备注
    add_textbox(slide, x + Inches(0.3), Inches(5.5), Inches(3.2), Inches(1.0),
                note, font_size=12, color=GRAY, alignment=PP_ALIGN.CENTER)

# ============================================================
# 第14页：总结回顾
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, LIGHT_BG)

add_shape_bg(slide, PRIMARY, Inches(0), Inches(0), Inches(0.15), prs.slide_height)

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
            '总结回顾', font_size=30, color=PRIMARY, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.0), Inches(2.5))

# 路径
add_textbox(slide, Inches(1.0), Inches(1.5), Inches(11), Inches(0.5),
            '一条清晰的路：', font_size=18, color=DARK, bold=True)

steps = [
    ('起点', '佛陀的选择：制度传承，以戒为师'),
    ('第一层', '规则大于个人 → 依法不依人'),
    ('第二层', '从"不敢"到"不忍" → 慈悲成为守法动力'),
    ('第三层', '法治与德治结合 → 完整约束到自觉的链条'),
    ('落脚', '知法守法 · 以信仰护法 · 服务社会'),
]

for i, (label, desc) in enumerate(steps):
    y = Inches(2.2) + Inches(i * 0.95)
    add_number_circle(slide, Inches(1.2), y, Inches(0.5), i+1, PRIMARY)
    add_textbox(slide, Inches(2.0), y + Inches(0.05), Inches(2.5), Inches(0.4),
                label, font_size=16, color=PRIMARY, bold=True)
    add_textbox(slide, Inches(4.5), y + Inches(0.05), Inches(7.5), Inches(0.4),
                desc, font_size=16, color=DARK)

# ============================================================
# 第15页：结语
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, RGBColor(0x1A, 0x0A, 0x0A))

add_shape_bg(slide, PRIMARY, Inches(0), Inches(0), prs.slide_width, Pt(6))

add_textbox(slide, Inches(1.5), Inches(1.0), Inches(10), Inches(0.6),
            '结语', font_size=24, color=ACCENT, alignment=PP_ALIGN.CENTER)

add_accent_line(slide, Inches(4.5), Inches(1.7), Inches(4.3), ACCENT, Pt(2))

add_textbox(slide, Inches(1.5), Inches(2.2), Inches(10), Inches(1.0),
            '让两束光交相辉映', font_size=36, color=WHITE, bold=True,
            alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(2), Inches(3.5), Inches(9.3), Inches(1.5),
            '信仰的光明  +  法治的阳光\n=  宗教中国化的前行之路', font_size=22, color=ACCENT,
            alignment=PP_ALIGN.CENTER)

add_accent_line(slide, Inches(3.5), Inches(5.2), Inches(6.3), ACCENT, Pt(1))

add_textbox(slide, Inches(2), Inches(5.5), Inches(9.3), Inches(1.0),
            '以戒为师，是信仰的定力\n依法护心，是时代的担当', font_size=24, color=WHITE, bold=True,
            alignment=PP_ALIGN.CENTER)

# ============================================================
# 第16页：封底
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, RGBColor(0x1A, 0x0A, 0x0A))

add_shape_bg(slide, PRIMARY, Inches(0), Inches(0), prs.slide_width, Pt(6))

add_textbox(slide, Inches(1.5), Inches(2.5), Inches(10), Inches(1.0),
            '谢谢', font_size=48, color=WHITE, bold=True,
            alignment=PP_ALIGN.CENTER)

add_accent_line(slide, Inches(4.5), Inches(3.6), Inches(4.3), ACCENT, Pt(2))

add_textbox(slide, Inches(1.5), Inches(4.0), Inches(10), Inches(0.8),
            '请各位老师、同学批评指正', font_size=20, color=ACCENT,
            alignment=PP_ALIGN.CENTER)

# 保存
output = '/home/admin/.openclaw/workspace/以戒为师依法护心-法治核心价值观演讲.pptx'
prs.save(output)
print(f'✅ PPT 已生成: {output}')
print(f'   共 {len(prs.slides)} 页')
