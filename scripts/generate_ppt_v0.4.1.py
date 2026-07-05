#!/usr/bin/env python3
"""
生成小龙虾网络 v0.4.1 项目汇报 PPT
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import sys

# 配色方案
COLOR_PRIMARY = RGBColor(0xFF, 0x6B, 0x35)    # 龙虾橙
COLOR_SECONDARY = RGBColor(0x00, 0x70, 0xC0)   # 科技蓝
COLOR_ACCENT = RGBColor(0x00, 0xB0, 0x50)      # 成功绿
COLOR_DARK = RGBColor(0x33, 0x33, 0x33)
COLOR_LIGHT = RGBColor(0xF5, 0xF5, 0xF5)
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)

def set_cell_text(cell, text, font_size=11, bold=False, color=COLOR_DARK):
    """设置单元格文本"""
    cell.text = ""
    p = cell.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE

def add_text_box(slide, left, top, width, height, text, font_size=18, 
                 bold=False, color=COLOR_DARK, alignment=PP_ALIGN.LEFT, font_name="微软雅黑"):
    """添加文本框"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    p.font.name = font_name
    return txBox

def add_bullet_list(slide, left, top, width, height, items, font_size=14, 
                    color=COLOR_DARK, font_name="微软雅黑"):
    """添加项目符号列表"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = Pt(8)
    
    return txBox

def add_shape_with_text(slide, shape_type, left, top, width, height, text, 
                        font_size=14, fill_color=None, text_color=COLOR_WHITE, bold=False):
    """添加带文字的图形"""
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    shape.text_frame.paragraphs[0].text = text
    shape.text_frame.paragraphs[0].font.size = Pt(font_size)
    shape.text_frame.paragraphs[0].font.color.rgb = text_color
    shape.text_frame.paragraphs[0].font.bold = bold
    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    return shape

def create_architecture_ppt():
    """创建架构汇报 PPT"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # ==================== 第1页：封面 ====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
    # 背景
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), 
                                 Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLOR_PRIMARY
    bg.line.fill.background()
    
    add_text_box(slide, Inches(1), Inches(2), Inches(11), Inches(1.5),
                 "小龙虾网络架构汇报", font_size=44, bold=True, color=COLOR_WHITE,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1), Inches(3.5), Inches(11), Inches(0.8),
                 "对话即创造：一人一世界的世界观", font_size=24, color=COLOR_WHITE,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1), Inches(4.5), Inches(11), Inches(0.6),
                 "一人一世界 × 世界是对话 × 世界是编程的", font_size=18, 
                 color=RGBColor(0xFF, 0xE0, 0xB0), alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.6),
                 "v0.4.1 整合版  |  信电大虾 & 诸葛斌  |  2026-06-24", 
                 font_size=16, color=RGBColor(0xFF, 0xE0, 0xB0), alignment=PP_ALIGN.CENTER)
    
    # ==================== 第2页：三层命题 ====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
                 "一、三层命题：完整的认知-生成-交互模型", font_size=32, 
                 bold=True, color=COLOR_PRIMARY)
    
    # 三个命题卡片
    propositions = [
        ("一人一世界", "认知编译系统", "每人一个种子参数", COLOR_PRIMARY),
        ("世界是对话", "认知张成与涌现", "对话输入新参数", COLOR_SECONDARY),
        ("世界是编程的", "过程生成引擎", "走到哪算到哪", COLOR_ACCENT),
    ]
    
    for i, (title, subtitle, desc, color) in enumerate(propositions):
        left = Inches(0.5 + i * 4.2)
        # 卡片背景
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.5), 
                                       Inches(3.8), Inches(2.5))
        card.fill.solid()
        card.fill.fore_color.rgb = color
        card.line.fill.background()
        
        add_text_box(slide, left + Inches(0.3), Inches(1.7), Inches(3.2), Inches(0.6),
                     title, font_size=28, bold=True, color=COLOR_WHITE, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, left + Inches(0.3), Inches(2.4), Inches(3.2), Inches(0.5),
                     subtitle, font_size=16, color=RGBColor(0xFF, 0xFF, 0xFF), alignment=PP_ALIGN.CENTER)
        add_text_box(slide, left + Inches(0.3), Inches(3.0), Inches(3.2), Inches(0.5),
                     desc, font_size=14, color=RGBColor(0xFF, 0xFF, 0xFF), alignment=PP_ALIGN.CENTER)
    
    # 串联逻辑
    add_text_box(slide, Inches(0.5), Inches(4.5), Inches(12), Inches(0.6),
                 "串联逻辑：种子 → 实时渲染 → 对话触发 → 涌现新地图", 
                 font_size=18, bold=True, color=COLOR_DARK)
    
    # 游戏实例
    items = [
        "🎮 玩家实例 = 认知编译系统（一人一世界）",
        "🌐 多人联机 = 交叉编译触发新事件（世界是对话）",
        "⚙️ 引擎底层 = 过程生成走到哪算到哪（世界是编程的）"
    ]
    add_bullet_list(slide, Inches(0.5), Inches(5.2), Inches(12), Inches(2), items, font_size=16)
    
    # ==================== 第3页：一人一世界 ====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
                 "二、一人一世界：没有"客观世界"这回事", font_size=32, 
                 bold=True, color=COLOR_PRIMARY)
    
    # 左侧：认知编译系统
    add_text_box(slide, Inches(0.5), Inches(1.2), Inches(6), Inches(0.5),
                 "认知编译系统", font_size=20, bold=True, color=COLOR_SECONDARY)
    
    tree_example = [
        "两个人站在同一棵树下：",
        "• 植物学家看到物种分类",
        "• 诗人看到生命隐喻",
        "• 樵夫看到木材",
        "",
        "树没变，编译结果不同",
        "",
        "感知系统 + 知识结构 + 价值坐标系",
        "共同编译出独属的"世界版本"",
        "",
        "你无法跳出自己的识去验证"客观""
    ]
    add_bullet_list(slide, Inches(0.5), Inches(1.8), Inches(6), Inches(3), tree_example, font_size=14)
    
    # 右侧：理论支撑
    add_text_box(slide, Inches(7), Inches(1.2), Inches(6), Inches(0.5),
                 "理论支撑", font_size=20, bold=True, color=COLOR_SECONDARY)
    
    theories = [
        "📖 佛教唯识学",
        "《成唯识论》：万法唯识",
        "一人一识，一识一界",
        "",
        "🧠 认知科学",
        "注意力机制决定什么进入意识",
        "没注意到的 = 没加载",
        "",
        "⚛️ 量子力学",
        "波函数不坍缩 = 没有确定位置",
        "观测就是渲染指令"
    ]
    add_bullet_list(slide, Inches(7), Inches(1.8), Inches(6), Inches(3), theories, font_size=14)
    
    # 结论
    add_text_box(slide, Inches(0.5), Inches(5.5), Inches(12), Inches(0.6),
                 "结论：世界不是静态数据库，是实时渲染引擎", 
                 font_size=20, bold=True, color=COLOR_ACCENT)
    
    # ==================== 第4页：世界是对话 ====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
                 "三、世界是对话：认知张成与涌现", font_size=32, 
                 bold=True, color=COLOR_PRIMARY)
    
    # 左侧：传统 vs 新理解
    add_text_box(slide, Inches(0.5), Inches(1.2), Inches(5.5), Inches(0.5),
                 "传统理解 vs 新理解", font_size=20, bold=True, color=COLOR_SECONDARY)
    
    traditional = [
        "传统理解：",
        "• A把已知信息传给B",
        "• 结果：各自知道了对方的信息",
        "",
        "新理解：",
        "• 两个认知系统交叉编译",
        "• 输出不属于任何单一主体的新结构",
        "",
        "例：同一道死活题，三种解法碰撞",
        "→ 第四种解法谁都没单独想到",
        "→ 这就是"新世界"的一个坐标点"
    ]
    add_bullet_list(slide, Inches(0.5), Inches(1.8), Inches(5.5), Inches(3.5), traditional, font_size=14)
    
    # 右侧：线性代数类比
    add_text_box(slide, Inches(6.5), Inches(1.2), Inches(6.5), Inches(0.5),
                 "线性代数类比", font_size=20, bold=True, color=COLOR_SECONDARY)
    
    analogy = [
        "一人一世界  →  基向量v₁,v₂,v₃  →  认知参数",
        "对    话    →  span{v₁,v₂,v₃}  →  认知张成",
        "新 世 界    →  高维子空间      →  涌现新解",
        "",
        "一人一世界 ≠ 孤岛，是乘法",
        "基向量越多、越不正交 → 张成空间越大"
    ]
    add_bullet_list(slide, Inches(6.5), Inches(1.8), Inches(6.5), Inches(2.5), analogy, font_size=14)
    
    # 华严经
    add_text_box(slide, Inches(0.5), Inches(5.2), Inches(12), Inches(0.8),
                 "华严经：因陀罗网 — 每颗宝珠映照所有宝珠，一即一切，一切即一", 
                 font_size=16, bold=True, color=COLOR_ACCENT)
    
    add_text_box(slide, Inches(0.5), Inches(6.0), Inches(12), Inches(0.5),
                 "结论：对话是输入新参数，让原来不存在的地图块被计算出来", 
                 font_size=18, bold=True, color=COLOR_PRIMARY)
    
    # ==================== 第5页：五层架构 ====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
                 "四、五层架构（v0.4.1 新增可靠通信层）", font_size=32, 
                 bold=True, color=COLOR_PRIMARY)
    
    layers = [
        ("应用层", "围棋训练 | 海报设计", COLOR_PRIMARY),
        ("运营层", "调度器 | Agent | 教练 | 监控", COLOR_SECONDARY),
        ("可靠通信层", "注册中心 | 可靠消息 | 故障切换", COLOR_ACCENT),  # 新增
        ("框架层", "节点 | 对话引擎 | 涌现检测 | 世界状态", COLOR_SECONDARY),
        ("基础设施层", "SSH通道 | OADP协议 | 配置 | 日志", COLOR_DARK),
    ]
    
    for i, (name, desc, color) in enumerate(layers):
        top = Inches(1.2 + i * 1.1)
        # 层背景
        layer = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), top, 
                                        Inches(11), Inches(0.9))
        layer.fill.solid()
        layer.fill.fore_color.rgb = color
        layer.line.fill.background()
        
        add_text_box(slide, Inches(1.2), top + Inches(0.1), Inches(3), Inches(0.7),
                     name, font_size=20, bold=True, color=COLOR_WHITE, alignment=PP_ALIGN.LEFT)
        add_text_box(slide, Inches(4.5), top + Inches(0.1), Inches(7), Inches(0.7),
                     desc, font_size=16, color=COLOR_WHITE, alignment=PP_ALIGN.LEFT)
        
        # 箭头
        if i < len(layers) - 1:
            arrow = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(6.5), top + Inches(0.9), 
                                            Inches(0.3), Inches(0.2))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = COLOR_DARK
            arrow.line.fill.background()
    
    # 可靠通信层标注
    add_text_box(slide, Inches(1), Inches(6.5), Inches(11), Inches(0.5),
                 "🆕 v0.4.1 新增：节点注册中心 + 可靠消息传递 + 多通道故障切换（NFS→SSH→HTTP→File）", 
                 font_size=14, bold=True, color=COLOR_ACCENT)
    
    # ==================== 第6页：v0.4.1 新增功能 ====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
                 "五、v0.4.1 新增功能", font_size=32, bold=True, color=COLOR_PRIMARY)
    
    # 节点注册中心
    add_text_box(slide, Inches(0.5), Inches(1.2), Inches(6), Inches(0.5),
                 "节点注册中心（NodeRegistry）", font_size=20, bold=True, color=COLOR_SECONDARY)
    
    registry_features = [
        "• 节点注册/注销：含能力声明和传输通道配置",
        "• 心跳检测：定期心跳，自动检测节点存活",
        "• 健康检查：全量检查，自动标记 offline/suspected",
        "• 节点发现：按类型/状态/能力查找",
        "• 持久化：JSON 文件存储，重启自动恢复",
        "• 回调机制：心跳回调、状态变化回调"
    ]
    add_bullet_list(slide, Inches(0.5), Inches(1.8), Inches(5.5), Inches(3), registry_features, font_size=13)
    
    # 可靠消息传递
    add_text_box(slide, Inches(6.5), Inches(1.2), Inches(6), Inches(0.5),
                 "可靠消息传递（Messenger）", font_size=20, bold=True, color=COLOR_SECONDARY)
    
    messenger_features = [
        "• 消息确认：ACK/NACK 机制",
        "• 自动重试：指数退避，可配置最大重试",
        "• 多通道故障切换：NFS→SSH→HTTP→File",
        "• 消息持久化：按状态分类存储",
        "• 消息过期：TTL 机制，自动清理",
        "• 优先级队列：支持消息优先级"
    ]
    add_bullet_list(slide, Inches(6.5), Inches(1.8), Inches(6), Inches(3), messenger_features, font_size=13)
    
    # 传输通道优先级表
    add_text_box(slide, Inches(0.5), Inches(4.8), Inches(12), Inches(0.5),
                 "传输通道优先级", font_size=20, bold=True, color=COLOR_SECONDARY)
    
    channels = [
        ("NFS", "1", "目录可写检查", "< 100ms"),
        ("SSH", "2", "连接测试", "< 500ms"),
        ("HTTP", "3", "HTTP 状态码", "< 1s"),
        ("File", "99", "目录可写", "< 50ms"),
    ]
    
    table = slide.shapes.add_table(len(channels)+1, 4, Inches(0.5), Inches(5.3), 
                                    Inches(12), Inches(1.5)).table
    table.columns[0].width = Inches(2)
    table.columns[1].width = Inches(2)
    table.columns[2].width = Inches(4)
    table.columns[3].width = Inches(4)
    
    headers = ["通道类型", "优先级", "故障检测", "切换时间"]
    for i, h in enumerate(headers):
        set_cell_text(table.rows[0].cells[i], h, font_size=12, bold=True, color=COLOR_WHITE)
        table.rows[0].cells[i].fill.solid()
        table.rows[0].cells[i].fill.fore_color.rgb = COLOR_PRIMARY
    
    for r, (ch, pri, det, time) in enumerate(channels):
        color = COLOR_ACCENT if r == 0 else COLOR_DARK
        for c, text in enumerate([ch, pri, det, time]):
            set_cell_text(table.rows[r+1].cells[c], text, font_size=12, color=color)
    
    # ==================== 第7页：理念与工程映射 ====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
                 "六、理念与工程映射", font_size=32, bold=True, color=COLOR_PRIMARY)
    
    add_text_box(slide, Inches(0.5), Inches(1.0), Inches(12), Inches(0.5),
                 "哲学命题 → 工程实现 → 验证指标", font_size=16, color=COLOR_SECONDARY)
    
    mappings = [
        ("一人一世界", "NodeRegistry 独立注册/心跳/状态", "节点存活率 > 99%"),
        ("世界是对话", "Messenger + ACK/NACK + 持久化", "消息投递成功率 > 99.5%"),
        ("世界是编程的", "多通道自动故障切换", "切换延迟 < 500ms"),
        ("因陀罗网", "全互联拓扑 + 健康检查", "节点发现 < 1s"),
        ("宝藏渲染", "涌现值计算 + 阈值触发", "涌现值 > 0.8 触发"),
    ]
    
    table = slide.shapes.add_table(len(mappings)+1, 3, Inches(0.5), Inches(1.5), 
                                    Inches(12), Inches(3.5)).table
    table.columns[0].width = Inches(3.5)
    table.columns[1].width = Inches(5)
    table.columns[2].width = Inches(3.5)
    
    headers = ["哲学命题", "工程实现（v0.4.1）", "验证指标"]
    for i, h in enumerate(headers):
        set_cell_text(table.rows[0].cells[i], h, font_size=14, bold=True, color=COLOR_WHITE)
        table.rows[0].cells[i].fill.solid()
        table.rows[0].cells[i].fill.fore_color.rgb = COLOR_PRIMARY
    
    for r, (phil, eng, metric) in enumerate(mappings):
        for c, text in enumerate([phil, eng, metric]):
            color = COLOR_ACCENT if c == 2 else COLOR_DARK
            set_cell_text(table.rows[r+1].cells[c], text, font_size=13, color=color)
    
    # 涌现实例
    add_text_box(slide, Inches(0.5), Inches(5.2), Inches(12), Inches(0.5),
                 "🌟 v0.4.1 涌现实例：虾尔(注册中心) + 诸葛马(SSH通道) → 可靠通信底座", 
                 font_size=16, bold=True, color=COLOR_ACCENT)
    
    # ==================== 第8页：测试与数据 ====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
                 "七、测试与数据", font_size=32, bold=True, color=COLOR_PRIMARY)
    
    # 测试统计
    add_text_box(slide, Inches(0.5), Inches(1.2), Inches(6), Inches(0.5),
                 "测试覆盖", font_size=20, bold=True, color=COLOR_SECONDARY)
    
    test_items = [
        "• 虾尔版：37 个测试全部通过",
        "• 诸葛马版：25 个测试全部通过",
        "• 总计：62 个测试全部通过 ✅",
        "",
        "• 注册中心：20 个测试",
        "• 可靠消息：12 个测试",
        "• 集成测试：5 个测试",
        "• 容错测试：2 个测试"
    ]
    add_bullet_list(slide, Inches(0.5), Inches(1.8), Inches(5.5), Inches(3), test_items, font_size=14)
    
    # 项目数据
    add_text_box(slide, Inches(6.5), Inches(1.2), Inches(6), Inches(0.5),
                 "项目数据", font_size=20, bold=True, color=COLOR_SECONDARY)
    
    data_items = [
        "• Python 文件：51+ 个",
        "• Markdown 文档：36+ 个",
        "• 总对局数：17,205+",
        "• 胜率：86%",
        "• 套利维度：5 维",
        "• 迭代速度：3 天 3 版本",
        "• 传输通道：4 种"
    ]
    add_bullet_list(slide, Inches(6.5), Inches(1.8), Inches(6), Inches(3), data_items, font_size=14)
    
    # ==================== 第9页：路线图 ====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
                 "八、开发路线图", font_size=32, bold=True, color=COLOR_PRIMARY)
    
    roadmap = [
        ("v0.4.0/v0.4.1", "注册中心 + 可靠消息 + SSH v2 + 部署脚本", "✅ 已完成", COLOR_ACCENT),
        ("v0.4.2", "SHA256 签名 + Prometheus 监控 + 告警系统", "🔄 规划中", RGBColor(0xFF, 0xC0, 0x00)),
        ("v0.5.0", "分布式注册中心 + 消息队列集群 + 跨域协作", "📋 计划中", COLOR_SECONDARY),
        ("v1.0.0", "生产级发布 + 多节点生态 + 完整文档", "🎯 目标", COLOR_PRIMARY),
    ]
    
    table = slide.shapes.add_table(len(roadmap)+1, 3, Inches(0.5), Inches(1.2), 
                                    Inches(12), Inches(3)).table
    table.columns[0].width = Inches(3)
    table.columns[1].width = Inches(6)
    table.columns[2].width = Inches(3)
    
    headers = ["版本", "核心交付", "状态"]
    for i, h in enumerate(headers):
        set_cell_text(table.rows[0].cells[i], h, font_size=14, bold=True, color=COLOR_WHITE)
        table.rows[0].cells[i].fill.solid()
        table.rows[0].cells[i].fill.fore_color.rgb = COLOR_PRIMARY
    
    for r, (ver, desc, status, color) in enumerate(roadmap):
        set_cell_text(table.rows[r+1].cells[0], ver, font_size=13, bold=True)
        set_cell_text(table.rows[r+1].cells[1], desc, font_size=13)
        set_cell_text(table.rows[r+1].cells[2], status, font_size=13, bold=True, color=color)
    
    # 下一步行动
    add_text_box(slide, Inches(0.5), Inches(4.5), Inches(12), Inches(0.5),
                 "下一步行动", font_size=20, bold=True, color=COLOR_SECONDARY)
    
    actions = [
        "☐ 与诸葛马确认对接方案",
        "☐ 建立 SSH 通信通道",
        "☐ 设计第一个对话触发器",
        "☐ 记录第一次对话的涌现值",
        "☐ 根据涌现结果调整网络参数"
    ]
    add_bullet_list(slide, Inches(0.5), Inches(5.0), Inches(12), Inches(2), actions, font_size=14)
    
    # ==================== 第10页：结论 ====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # 背景
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), 
                                 Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLOR_PRIMARY
    bg.line.fill.background()
    
    add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(1),
                 "九、结论", font_size=40, bold=True, color=COLOR_WHITE,
                 alignment=PP_ALIGN.CENTER)
    
    conclusions = [
        "一人一世界 ≠ 各过各的孤岛",
        "而是：承认差异是常态，理解是努力的结果，共识是动态的临时协议",
        "",
        "对话不是"交换已知信息"",
        "是输入新的参数到生成引擎，让原来不存在的地图块被计算出来",
        "",
        "小龙虾网络就是这个引擎的工程实现",
        "",
        "你不停对话，世界就不停扩展 🦞⚡️"
    ]
    
    y = Inches(2.8)
    for text in conclusions:
        size = 24 if text else 16
        bold = True if "你不停对话" in text else False
        color = COLOR_WHITE if text else COLOR_WHITE
        add_text_box(slide, Inches(1), y, Inches(11), Inches(0.5),
                     text, font_size=size, bold=bold, color=color, alignment=PP_ALIGN.CENTER)
        y += Inches(0.5)
    
    # 保存
    output_path = '/tmp/小龙虾网络架构汇报-v0.4.1.pptx'
    prs.save(output_path)
    print(f"✅ 架构汇报 PPT 已生成: {output_path}")
    return output_path

if __name__ == "__main__":
    create_architecture_ppt()
