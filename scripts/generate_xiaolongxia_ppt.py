#!/usr/bin/env python3
"""
生成「小龙虾智能体落地推广」科普汇报PPT
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ====== 配色方案 ======
BG_DARK = RGBColor(0x1A, 0x1A, 0x2E)       # 深蓝黑背景
BG_CARD = RGBColor(0x16, 0x21, 0x3E)        # 卡片背景
ACCENT_ORANGE = RGBColor(0xFF, 0x6B, 0x35)  # 主色（活力橙）
ACCENT_BLUE = RGBColor(0x4D, 0xAB, 0xF7)    # 辅助蓝
ACCENT_GREEN = RGBColor(0x2E, 0xCC, 0x71)   # 绿色
ACCENT_RED = RGBColor(0xE7, 0x4C, 0x3C)     # 红色
ACCENT_PURPLE = RGBColor(0x9B, 0x59, 0xB6)  # 紫色
ACCENT_YELLOW = RGBColor(0xF1, 0xC4, 0x0F)  # 黄色
TEXT_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_LIGHT = RGBColor(0xCC, 0xCC, 0xCC)
TEXT_DIM = RGBColor(0x99, 0x99, 0xAA)
LINE_ORANGE = RGBColor(0xFF, 0x8C, 0x5A)

def add_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape_bg(slide, left, top, width, height, color, alpha=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if alpha is not None:
        from lxml import etree
        ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
        solidFill = shape.fill._fill
        srgb = solidFill.find(f'.//{{{ns}}}srgbClr')
        if srgb is not None:
            alpha_elem = etree.SubElement(srgb, f'{{{ns}}}alpha')
            alpha_elem.set('val', str(int(alpha * 1000)))
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=TEXT_WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name='Microsoft YaHei'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_card(slide, left, top, width, height, bg_color=BG_CARD):
    card = add_shape_bg(slide, left, top, width, height, bg_color)
    return card

def add_icon_text(slide, left, top, icon_emoji, title, body, title_size=22, body_size=16,
                  title_color=ACCENT_ORANGE, body_color=TEXT_LIGHT):
    add_text_box(slide, left, top, Inches(0.6), Inches(0.5), icon_emoji,
                 font_size=title_size, bold=True, color=title_color)
    add_text_box(slide, left + Inches(0.6), top, Inches(5), Inches(0.5), title,
                 font_size=title_size, bold=True, color=TEXT_WHITE)
    add_text_box(slide, left + Inches(0.6), top + Inches(0.5), Inches(5.5), Inches(1.5), body,
                 font_size=body_size, color=body_color)

def add_two_column_layout(slide):
    """在内容区画一条分隔线"""
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(6.7), Inches(2.2), Inches(0.03), Inches(4.5))
    line.fill.solid()
    line.fill.fore_color.rgb = LINE_ORANGE
    line.line.fill.background()

def set_tf(tf, paragraphs, default_size=16, default_color=TEXT_LIGHT):
    """设置text_frame的多行段落"""
    for i, (text, size, color, bold) in enumerate(paragraphs):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size or default_size)
        p.font.color.rgb = color or default_color
        p.font.bold = bold if bold else False
        p.font.name = 'Microsoft YaHei'
        p.space_after = Pt(4)

# ============================================================
# SLIDE 1: 封面
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide)

# 装饰色块
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), ACCENT_ORANGE)

# 大标题
add_text_box(slide, Inches(1.5), Inches(1.5), Inches(10), Inches(1.2),
             '🦞 小龙虾智能体', font_size=52, bold=True, color=ACCENT_ORANGE,
             alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1.5), Inches(2.8), Inches(10), Inches(1.2),
             '落地推广之路', font_size=44, bold=True, color=TEXT_WHITE,
             alignment=PP_ALIGN.CENTER)

# 副标题
add_text_box(slide, Inches(2), Inches(4.2), Inches(9.3), Inches(0.8),
             '从互联网+到人工智能+ —— 一个AI智能体的诞生与成长',
             font_size=24, color=TEXT_DIM, alignment=PP_ALIGN.CENTER)

# 底部信息
add_text_box(slide, Inches(2), Inches(5.8), Inches(9.3), Inches(0.5),
             '诸葛斌  |  2026年6月',
             font_size=18, color=TEXT_DIM, alignment=PP_ALIGN.CENTER)

# 底部装饰线
add_shape_bg(slide, Inches(4), Inches(5.5), Inches(5.3), Inches(0.04), ACCENT_ORANGE)

# ============================================================
# SLIDE 2: 目录
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), ACCENT_ORANGE)

add_text_box(slide, Inches(1), Inches(0.5), Inches(11), Inches(0.8),
             '📋 汇报目录', font_size=36, bold=True, color=ACCENT_ORANGE)

toc_items = [
    ('01', '时代背景', '从互联网+到人工智能+的十年跨越'),
    ('02', '历史启示', '互联网+创业潮的经验与教训'),
    ('03', '小龙虾是谁', '智能体的能力矩阵与定位'),
    ('04', '推广策略', '三阶段落地路径'),
    ('05', '优势与挑战', 'AI时代的机遇与风险'),
    ('06', '未来展望', '从小工具到大生态'),
]

for i, (num, title, desc) in enumerate(toc_items):
    y = Inches(1.8) + Inches(i * 0.85)
    card = add_card(slide, Inches(1.5), y, Inches(10), Inches(0.7), BG_CARD)
    
    add_text_box(slide, Inches(1.8), y + Inches(0.08), Inches(0.8), Inches(0.5),
                 num, font_size=28, bold=True, color=ACCENT_ORANGE)
    add_text_box(slide, Inches(2.7), y + Inches(0.05), Inches(4), Inches(0.5),
                 title, font_size=22, bold=True, color=TEXT_WHITE)
    add_text_box(slide, Inches(2.7), y + Inches(0.4), Inches(8), Inches(0.3),
                 desc, font_size=14, color=TEXT_DIM)

# ============================================================
# SLIDE 3: 时代背景
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), ACCENT_ORANGE)

add_text_box(slide, Inches(1), Inches(0.5), Inches(11), Inches(0.8),
             '🕐 时代背景：从互联网+ 到人工智能+',
             font_size=36, bold=True, color=ACCENT_ORANGE)

# 左：互联网+
add_card(slide, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.2), BG_CARD)
add_text_box(slide, Inches(1.1), Inches(1.8), Inches(5), Inches(0.6),
             '🌐 互联网+ 时代（2015）', font_size=26, bold=True, color=ACCENT_BLUE)

left_items = [
    ('📄 政策文件', '国发〔2015〕40号《关于积极推进"互联网+"行动的指导意见》'),
    ('🎯 核心逻辑', '连接与赋能 — 用互联网连接传统产业'),
    ('🔑 关键词', '融合、创业创新、平台经济、O2O'),
    ('🏗️ 11个重点领域', '创业创新、协同制造、电商、物流、金融、养老...'),
    ('📱 技术基础', '移动互联网、云计算、大数据、IoT'),
    ('🌍 国际环境', '中国消费互联网崛起，全球领先'),
]
for i, (label, val) in enumerate(left_items):
    y = Inches(2.6) + Inches(i * 0.65)
    add_text_box(slide, Inches(1.3), y, Inches(1.8), Inches(0.3),
                 label, font_size=15, bold=True, color=ACCENT_BLUE)
    add_text_box(slide, Inches(3.1), y, Inches(3.1), Inches(0.55),
                 val, font_size=14, color=TEXT_LIGHT)

# 右：人工智能+
add_card(slide, Inches(7), Inches(1.6), Inches(5.6), Inches(5.2), BG_CARD)
add_text_box(slide, Inches(7.3), Inches(1.8), Inches(5), Inches(0.6),
             '🤖 人工智能+ 时代（2025）', font_size=26, bold=True, color=ACCENT_ORANGE)

right_items = [
    ('📄 政策起点', '2024年政府工作报告首次提出"人工智能+"行动'),
    ('🎯 核心逻辑', '替代与重塑 — 用AI重构产业逻辑与生产方式'),
    ('🔑 关键词', '大模型、新质生产力、自主可控、智能体'),
    ('🏗️ 重点领域', '大模型研发、AI+制造/医疗/教育、算力基建'),
    ('🧠 技术基础', '大语言模型、生成式AI、AGI竞争'),
    ('🌍 国际环境', '中美AI博弈、芯片管制、技术追赶'),
]
for i, (label, val) in enumerate(right_items):
    y = Inches(2.6) + Inches(i * 0.65)
    add_text_box(slide, Inches(7.5), y, Inches(1.8), Inches(0.3),
                 label, font_size=15, bold=True, color=ACCENT_ORANGE)
    add_text_box(slide, Inches(9.3), y, Inches(3.1), Inches(0.55),
                 val, font_size=14, color=TEXT_LIGHT)

# 底部总结
add_text_box(slide, Inches(1), Inches(6.9), Inches(11), Inches(0.5),
             '💡 一句话：互联网+ = 连接万物  →  人工智能+ = 理解万物',
             font_size=18, bold=True, color=ACCENT_YELLOW, alignment=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 4: 历史启示
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), ACCENT_ORANGE)

add_text_box(slide, Inches(1), Inches(0.5), Inches(11), Inches(0.8),
             '📖 历史启示：互联网+创业潮的经验与教训',
             font_size=36, bold=True, color=ACCENT_ORANGE)

# 案例卡片
cases = [
    ('🔶 王兴：连环创业',
     '校内网 → 饭否网 → 美团\n四次创业三次失败，千团大战中活下来\n核心：快速试错，抄作业本土化',
     ACCENT_ORANGE),
    ('🔷 张一鸣：算法重构',
     '2012年创办今日头条，押注"算法推荐"\n所有人质疑"机器懂什么内容？"\n核心：用技术解决老问题，不怕"非共识"',
     ACCENT_BLUE),
    ('🔴 千团大战',
     '5000多家团购网站 → 最终只剩几家\n烧钱获客 → 99%灰飞烟灭\n核心：风口论 + 资本驱动 ≠ 长期竞争力',
     ACCENT_RED),
]

for i, (title, body, color) in enumerate(cases):
    x = Inches(0.8) + Inches(i * 4.1)
    add_card(slide, x, Inches(1.6), Inches(3.8), Inches(2.5), BG_CARD)
    add_text_box(slide, x + Inches(0.3), Inches(1.8), Inches(3.2), Inches(0.5),
                 title, font_size=20, bold=True, color=color)
    add_text_box(slide, x + Inches(0.3), Inches(2.4), Inches(3.2), Inches(1.5),
                 body, font_size=14, color=TEXT_LIGHT)

# 教训总结
add_card(slide, Inches(0.8), Inches(4.4), Inches(11.8), Inches(2.8), BG_CARD)
add_text_box(slide, Inches(1.2), Inches(4.5), Inches(11), Inches(0.5),
             '⚠️ 互联网+时代的三大教训', font_size=24, bold=True, color=ACCENT_YELLOW)

lessons = [
    '❌ 教训一：不要做"工具"，要做"解决方案" — 用户不在乎技术多强，只在乎问题能不能解决',
    '❌ 教训二：找准种子用户，别一上来就搞"平台" — 先在一个场景打透，再复制扩张',
    '❌ 教训三：技术壁垒 ≠ 商业壁垒 — 真正的壁垒是场景理解、数据积累和用户信任',
]
for i, lesson in enumerate(lessons):
    add_text_box(slide, Inches(1.5), Inches(5.2) + Inches(i * 0.55), Inches(10.5), Inches(0.5),
                 lesson, font_size=16, color=TEXT_LIGHT)

# ============================================================
# SLIDE 5: 两个时代的相似与不同
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), ACCENT_ORANGE)

add_text_box(slide, Inches(1), Inches(0.5), Inches(11), Inches(0.8),
             '🔍 两个时代的相似与不同',
             font_size=36, bold=True, color=ACCENT_ORANGE)

# 相似之处
add_card(slide, Inches(0.8), Inches(1.6), Inches(5.8), Inches(5.2), BG_CARD)
add_text_box(slide, Inches(1.1), Inches(1.8), Inches(5), Inches(0.5),
             '✅ 相似之处', font_size=26, bold=True, color=ACCENT_GREEN)

sims = [
    '📢 "全民狂欢"的叙事结构\n万物皆可互联网 → 万物皆可AI',
    '⚔️ "千团大战" → "百模大战"\n5000家团购 vs 数十个大模型\n都用烧钱/烧算力抢用户',
    '⏰ "工具红利期"的时间窗口\n先上车的吃肉，后上车的喝汤\n但窗口期正在快速缩短',
    '🫧 泡沫与价值的辩证关系\n泡沫是技术扩散的必要成本\n99%会死，但1%会成长为巨头',
]
for i, sim in enumerate(sims):
    y = Inches(2.5) + Inches(i * 1.05)
    add_text_box(slide, Inches(1.3), y, Inches(5), Inches(0.9),
                 sim, font_size=15, color=TEXT_LIGHT)

# 不同之处
add_card(slide, Inches(7), Inches(1.6), Inches(5.8), Inches(5.2), BG_CARD)
add_text_box(slide, Inches(7.3), Inches(1.8), Inches(5), Inches(0.5),
             '🆕 不同之处', font_size=26, bold=True, color=ACCENT_ORANGE)

diffs = [
    '📉 创业门槛剧变\n互联网+需技术团队 → AI时代一人即可\n但"活得久"的门槛反而更高',
    '💰 资本逻辑变化\nVC"融资→烧钱→上市" → AI时代"烧算力=烧钱"\n初创公司面临"烧钱就烧死"的两难',
    '🔄 技术替代的深度\n互联网+解决信息不对称 → AI解决认知能力替代',
    '📜 政策环境差异\n互联网+先发展后治理 → AI时代发展与治理同步',
]
for i, diff in enumerate(diffs):
    y = Inches(2.5) + Inches(i * 1.05)
    add_text_box(slide, Inches(7.5), y, Inches(5), Inches(0.9),
                 diff, font_size=15, color=TEXT_LIGHT)

# ============================================================
# SLIDE 6: 小龙虾智能体是谁
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), ACCENT_ORANGE)

add_text_box(slide, Inches(1), Inches(0.5), Inches(11), Inches(0.8),
             '🦞 小龙虾智能体是谁？',
             font_size=36, bold=True, color=ACCENT_ORANGE)

# 简介
add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(0.6),
             '不是AI玩具，而是已经在跑真实业务的多智能体系统',
             font_size=22, bold=True, color=ACCENT_YELLOW, alignment=PP_ALIGN.CENTER)

# 实例卡片
add_card(slide, Inches(0.8), Inches(2.3), Inches(5.6), Inches(2.2), BG_CARD)
add_text_box(slide, Inches(1.1), Inches(2.4), Inches(5), Inches(0.5),
             '🦐 虾尔（小龙虾-诸葛虾）', font_size=22, bold=True, color=ACCENT_ORANGE)
add_text_box(slide, Inches(1.3), Inches(2.9), Inches(5), Inches(1.5),
             '• 主智能体，基于OpenClaw框架\n• 教学辅助、课表转换、论文指导\n• 钉钉集成、智能提醒、课程分析\n• 20+技能模块，支持多智能体协作',
             font_size=15, color=TEXT_LIGHT)

add_card(slide, Inches(7), Inches(2.3), Inches(5.6), Inches(2.2), BG_CARD)
add_text_box(slide, Inches(7.3), Inches(2.4), Inches(5), Inches(0.5),
             '🐴 诸葛马（Hermes）', font_size=22, bold=True, color=ACCENT_BLUE)
add_text_box(slide, Inches(7.5), Inches(2.9), Inches(5), Inches(1.5),
             '• 协作节点，独立服务器部署\n• NFS双向通信，能力共享\n• 任务协同，消息中转\n• 与虾尔构成分布式智能体网络',
             font_size=15, color=TEXT_LIGHT)

# 技术架构
add_card(slide, Inches(0.8), Inches(4.7), Inches(11.8), Inches(2.5), BG_CARD)
add_text_box(slide, Inches(1.1), Inches(4.8), Inches(11), Inches(0.5),
             '🏗️ 技术架构特点', font_size=22, bold=True, color=ACCENT_PURPLE)

arch_items = [
    ('🔗 分布式通信', 'NFS共享目录实现实例间消息传递，30分钟自动同步'),
    ('🧩 技能生态', 'SKILL.md标准化技能接口，支持热插拔扩展'),
    ('📱 多端触达', '钉钉、QQ Bot、Web等多渠道消息推送'),
    ('🔒 本地部署', '数据不出校，满足教育数据安全合规要求'),
]
for i, (label, desc) in enumerate(arch_items):
    x = Inches(1.3) + Inches(i * 2.9)
    add_text_box(slide, x, Inches(5.4), Inches(2.6), Inches(0.35),
                 label, font_size=16, bold=True, color=ACCENT_PURPLE)
    add_text_box(slide, x, Inches(5.8), Inches(2.6), Inches(0.8),
                 desc, font_size=13, color=TEXT_LIGHT)

# ============================================================
# SLIDE 7: 能力矩阵
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), ACCENT_ORANGE)

add_text_box(slide, Inches(1), Inches(0.5), Inches(11), Inches(0.8),
             '⚡ 已落地的能力矩阵',
             font_size=36, bold=True, color=ACCENT_ORANGE)

skills = [
    ('📅 课表转换', '中方→英方\n自动转换教学日历', ACCENT_ORANGE),
    ('📝 论文指导', '评审管理\n格式审查\n定稿确认', ACCENT_BLUE),
    ('📊 教学分析', '学生评价\n五维评分\n成效报告', ACCENT_GREEN),
    ('💬 钉钉集成', '消息推送\n案例导出\n群管理', ACCENT_PURPLE),
    ('⏰ 智能提醒', '定时任务\n周期性提醒\n多渠道送达', ACCENT_YELLOW),
    ('📚 学术写作', 'LaTeX支持\n文献检索\n引用管理', ACCENT_RED),
    ('🌤️ 生活助手', '天气查询\n网络搜索\n媒体处理', TEXT_LIGHT),
    ('🔗 多智能体', 'NFS通信\n任务协同\n能力共享', ACCENT_ORANGE),
]

for i, (title, desc, color) in enumerate(skills):
    col = i % 4
    row = i // 4
    x = Inches(0.8) + Inches(col * 3.1)
    y = Inches(1.6) + Inches(row * 2.7)
    add_card(slide, x, y, Inches(2.8), Inches(2.4), BG_CARD)
    # 左侧色条
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  x, y, Inches(0.08), Inches(2.4))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    
    add_text_box(slide, x + Inches(0.3), y + Inches(0.2), Inches(2.3), Inches(0.5),
                 title, font_size=20, bold=True, color=color)
    add_text_box(slide, x + Inches(0.3), y + Inches(0.8), Inches(2.3), Inches(1.4),
                 desc, font_size=14, color=TEXT_LIGHT)

# ============================================================
# SLIDE 8: 推广核心原则
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), ACCENT_ORANGE)

add_text_box(slide, Inches(1), Inches(0.5), Inches(11), Inches(0.8),
             '🎯 推广核心原则：从教训到行动',
             font_size=36, bold=True, color=ACCENT_ORANGE)

principles = [
    ('❌ 互联网+的错误', '✅ 小龙虾的做法', ''),
    ('"我们的系统很强大"\n到处推销技术能力', '"帮你把2小时课表变成3分钟"\n从痛点出发，不谈技术谈价值', '原则一：不做工具，做解决方案'),
    ('一上来就搞"平台"\n全国300城同时开战', '先吃透自己的场景\n5个深度用户 > 500个注册用户', '原则二：单场景打透，再复制'),
    ('以为"技术复制不了"就是护城河\n结果巨头三个月就抄出来', '真正的壁垒是：\n场景理解 + 数据积累 + 用户信任', '原则三：场景深度 > 技术壁垒'),
]

for i, (wrong, right, principle) in enumerate(principles):
    y = Inches(1.6) + Inches(i * 1.7)
    if i == 0:
        add_text_box(slide, Inches(0.8), y, Inches(5.6), Inches(0.5),
                     wrong, font_size=20, bold=True, color=ACCENT_RED)
        add_text_box(slide, Inches(7), y, Inches(5.6), Inches(0.5),
                     right, font_size=20, bold=True, color=ACCENT_GREEN)
    else:
        add_card(slide, Inches(0.8), y - Inches(0.1), Inches(5.6), Inches(1.5), BG_CARD)
        add_card(slide, Inches(7), y - Inches(0.1), Inches(5.6), Inches(1.5), BG_CARD)
        add_text_box(slide, Inches(1.1), y, Inches(5), Inches(1.3),
                     wrong, font_size=14, color=TEXT_DIM)
        add_text_box(slide, Inches(7.3), y, Inches(5), Inches(1.3),
                     right, font_size=14, color=TEXT_LIGHT)
        # 原则标签
        add_text_box(slide, Inches(5.2), y + Inches(0.4), Inches(3), Inches(0.4),
                     '👉 ' + principle, font_size=15, bold=True, color=ACCENT_YELLOW)

# ============================================================
# SLIDE 9: 推广路径
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), ACCENT_ORANGE)

add_text_box(slide, Inches(1), Inches(0.5), Inches(11), Inches(0.8),
             '🗺️ 三阶段推广路径',
             font_size=36, bold=True, color=ACCENT_ORANGE)

phases = [
    ('阶段一：验证期', '1-3个月', '证明"有人真的在用，且离不开"', ACCENT_BLUE, [
        '📌 核心指标：日活、使用频次、替代人工时间',
        '📌 完善3-5个核心技能的用户体验',
        '📌 记录每次使用节省的时间',
        '📌 收集用户原话反馈',
        '🎯 推广方式：口碑，不做营销',
    ]),
    ('阶段二：种子期', '3-6个月', '10-20个深度用户，形成"离不开"', ACCENT_ORANGE, [
        '📌 核心指标：留存率、功能深度、转介绍率',
        '📌 根据反馈迭代技能',
        '📌 整理标准化部署文档',
        '📌 建立用户交流群（钉钉）',
        '🎯 推广方式：同行推荐、教研会展示',
    ]),
    ('阶段三：扩张期', '6-12个月', '50+用户，覆盖多个院校/专业', ACCENT_GREEN, [
        '📌 核心指标：付费转化、NPS净推荐值',
        '📌 SaaS化 vs 本地部署路线选择',
        '📌 "一键安装"脚本降低门槛',
        '📌 建立技能市场（社区贡献）',
        '🎯 推广方式：行业会议、院校合作',
    ]),
]

for i, (title, period, goal, color, items) in enumerate(phases):
    x = Inches(0.8) + Inches(i * 4.1)
    add_card(slide, x, Inches(1.6), Inches(3.8), Inches(5.3), BG_CARD)
    # 顶部色条
    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                      x, Inches(1.6), Inches(3.8), Inches(0.06))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = color
    top_bar.line.fill.background()
    
    add_text_box(slide, x + Inches(0.3), Inches(1.8), Inches(3.2), Inches(0.4),
                 title, font_size=22, bold=True, color=color)
    add_text_box(slide, x + Inches(0.3), Inches(2.2), Inches(3.2), Inches(0.3),
                 period, font_size=16, color=TEXT_DIM)
    
    add_text_box(slide, x + Inches(0.3), Inches(2.6), Inches(3.2), Inches(0.5),
                 goal, font_size=14, bold=True, color=ACCENT_YELLOW)
    
    for j, item in enumerate(items):
        add_text_box(slide, x + Inches(0.3), Inches(3.3) + Inches(j * 0.65),
                     Inches(3.2), Inches(0.55),
                     item, font_size=13, color=TEXT_LIGHT)

# ============================================================
# SLIDE 10: AI时代的独特优势
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), ACCENT_ORANGE)

add_text_box(slide, Inches(1), Inches(0.5), Inches(11), Inches(0.8),
             '🚀 AI时代的独特优势（互联网+时代没有的）',
             font_size=36, bold=True, color=ACCENT_ORANGE)

advantages = [
    ('💰 冷启动成本极低',
     '互联网+时代：找CTO、租服务器、写代码，至少3个月\nAI时代：一个人 + OpenClaw + 现有技能 = 可用的AI产品\n启动成本几乎为零',
     ACCENT_GREEN),
    ('🎨 个性化的规模化',
     '互联网+时代：标准化才能规模化，个性化 = 高成本\nAI时代：同时服务100个用户，每个对话都是个性化的\n"私人定制"体验 + 大规模服务',
     ACCENT_BLUE),
    ('⚡ 指数级迭代速度',
     '互联网+时代：功能从需求到上线，最快一周\nAI时代：加一个新技能，几小时就能上线\n迭代速度是传统软件的10倍以上',
     ACCENT_ORANGE),
    ('🦸 "超级个体"模式',
     '互联网+时代：做产品需要产品、开发、运营、客服\nAI时代：你 + 智能体 = 完整团队\n智能体自动处理客服、运营、甚至部分开发',
     ACCENT_PURPLE),
]

for i, (title, desc, color) in enumerate(advantages):
    col = i % 2
    row = i // 2
    x = Inches(0.8) + Inches(col * 6.2)
    y = Inches(1.6) + Inches(row * 2.8)
    add_card(slide, x, y, Inches(5.9), Inches(2.5), BG_CARD)
    
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  x, y, Inches(0.08), Inches(2.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    
    add_text_box(slide, x + Inches(0.3), y + Inches(0.15), Inches(5.4), Inches(0.45),
                 title, font_size=22, bold=True, color=color)
    add_text_box(slide, x + Inches(0.3), y + Inches(0.7), Inches(5.4), Inches(1.6),
                 desc, font_size=14, color=TEXT_LIGHT)

# ============================================================
# SLIDE 11: 挑战与应对
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), ACCENT_ORANGE)

add_text_box(slide, Inches(1), Inches(0.5), Inches(11), Inches(0.8),
             '⚠️ 潜在风险与应对策略',
             font_size=36, bold=True, color=ACCENT_ORANGE)

risks = [
    ('🏢 巨头入场', '阿里/腾讯推出教育智能体平台', '守住垂直场景，做他们不愿做的"脏活累活"', ACCENT_RED),
    ('🔄 技术平权', '人人都能搭建智能体', '壁垒不在工具，在场景理解和数据积累', ACCENT_ORANGE),
    ('🔒 数据安全', '教学数据、学生信息合规', '优先本地部署方案，数据不出校', ACCENT_BLUE),
    ('👤 用户习惯', '教师对AI接受度参差不齐', '先服务"早期采用者"，让他们影响观望者', ACCENT_GREEN),
    ('🔗 依赖风险', '过度依赖底层平台/API', '保持架构灵活性，核心逻辑可迁移', ACCENT_PURPLE),
]

for i, (risk, desc, response, color) in enumerate(risks):
    y = Inches(1.6) + Inches(i * 1.05)
    add_card(slide, Inches(0.8), y, Inches(11.8), Inches(0.9), BG_CARD)
    
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(0.8), y, Inches(0.06), Inches(0.9))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    
    add_text_box(slide, Inches(1.1), y + Inches(0.05), Inches(2.5), Inches(0.35),
                 risk, font_size=18, bold=True, color=color)
    add_text_box(slide, Inches(3.8), y + Inches(0.05), Inches(3.5), Inches(0.35),
                 desc, font_size=14, color=TEXT_DIM)
    add_text_box(slide, Inches(7.5), y + Inches(0.05), Inches(5), Inches(0.7),
                 '✅ ' + response, font_size=14, color=TEXT_LIGHT)

# ============================================================
# SLIDE 12: 总结
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), ACCENT_ORANGE)

add_text_box(slide, Inches(1), Inches(0.5), Inches(11), Inches(0.8),
             '🎯 总结',
             font_size=36, bold=True, color=ACCENT_ORANGE)

# 核心结论卡片
add_card(slide, Inches(1.5), Inches(1.8), Inches(10.3), Inches(1.2), BG_CARD)
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                              Inches(1.5), Inches(1.8), Inches(0.08), Inches(1.2))
bar.fill.solid()
bar.fill.fore_color.rgb = ACCENT_ORANGE
bar.line.fill.background()

add_text_box(slide, Inches(2), Inches(1.9), Inches(9.5), Inches(1),
             '💡 "互联网+时代教我们"敢不敢做"，AI时代考验的是"做得有多深"。"',
             font_size=26, bold=True, color=ACCENT_YELLOW, alignment=PP_ALIGN.CENTER)

# 三个关键信息
key_points = [
    ('🦞 不要做"下一个钉钉"', '做中英合作办学领域最好用的AI助手，这个市场就足够大了。垂直深耕 > 横向扩张。', ACCENT_ORANGE),
    ('🌱 从自己出发，先做种子用户', '你既是开发者，也是第一个深度用户。把自己的体验打磨到极致，自然会有人跟随。', ACCENT_GREEN),
    ('⏰ 时间站在你这边', '技术平权让个人创业者的窗口变短，但"超级个体"的可能性前所未有。现在是最好的时机。', ACCENT_BLUE),
]

for i, (title, desc, color) in enumerate(key_points):
    y = Inches(3.4) + Inches(i * 1.2)
    add_card(slide, Inches(1.5), y, Inches(10.3), Inches(1), BG_CARD)
    
    add_text_box(slide, Inches(2), y + Inches(0.1), Inches(4), Inches(0.4),
                 title, font_size=20, bold=True, color=color)
    add_text_box(slide, Inches(2), y + Inches(0.5), Inches(9.5), Inches(0.4),
                 desc, font_size=15, color=TEXT_LIGHT)

# ============================================================
# SLIDE 13: 结尾页
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), ACCENT_ORANGE)
add_shape_bg(slide, Inches(0), Inches(7.42), Inches(13.333), Inches(0.08), ACCENT_ORANGE)

add_text_box(slide, Inches(2), Inches(2.2), Inches(9.3), Inches(1),
             '🦞 小龙虾智能体', font_size=48, bold=True, color=ACCENT_ORANGE,
             alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(2), Inches(3.3), Inches(9.3), Inches(1),
             '技术浪潮不会等人，但站在浪尖的人\n从来不是技术最强的，而是行动最快的',
             font_size=28, color=TEXT_WHITE, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(2), Inches(4.8), Inches(9.3), Inches(0.6),
             '感谢聆听', font_size=36, bold=True, color=ACCENT_YELLOW,
             alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(2), Inches(5.8), Inches(9.3), Inches(0.5),
             '诸葛斌  |  2026年6月',
             font_size=18, color=TEXT_DIM, alignment=PP_ALIGN.CENTER)

# 保存
output_path = '/home/admin/.openclaw/workspace/小龙虾智能体落地推广汇报.pptx'
prs.save(output_path)
print(f'✅ PPT已保存: {output_path}')
