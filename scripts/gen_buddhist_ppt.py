#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Generate academic PPT for 佛教文化与宗教中国化研究"""

from pptx import Presentation
from pptx.util import Inches, Pt, Cm, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Cm(33.867)
prs.slide_height = Cm(19.05)

# === Color Palette ===
BG_DARK = RGBColor(0x1B, 0x2A, 0x4A)      # 深蓝背景
BG_ACCENT = RGBColor(0x0F, 0x1F, 0x3D)     # 更深的蓝
GOLD = RGBColor(0xD4, 0xA5, 0x47)          # 金色（佛教色）
GOLD_LIGHT = RGBColor(0xE8, 0xC8, 0x72)    # 浅金
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
MID_GRAY = RGBColor(0x99, 0x99, 0x99)
TEXT_BODY = RGBColor(0xE0, 0xE0, 0xE0)
ACCENT_RED = RGBColor(0xC0, 0x39, 0x2B)
ACCENT_GREEN = RGBColor(0x27, 0xAE, 0x60)

def add_bg(slide, color=BG_DARK):
    """Fill slide background"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape_bg(slide, left, top, width, height, color, alpha=None):
    """Add a colored rectangle as background element"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_gold_line(slide, left, top, width=Cm(8), height=Pt(3)):
    """Add a decorative gold line"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = GOLD
    shape.line.fill.background()

def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name='微软雅黑', line_spacing=1.5):
    """Add a text box with formatted text"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    if line_spacing:
        p.space_after = Pt(font_size * 0.3)
    return txBox

def add_bullet_list(slide, left, top, width, height, items, font_size=16,
                    color=TEXT_BODY, bullet_color=GOLD, line_spacing=Pt(8)):
    """Add a bulleted list"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        # Bullet character
        run1 = p.add_run()
        run1.text = '● '
        run1.font.size = Pt(font_size - 2)
        run1.font.color.rgb = bullet_color
        run1.font.name = '微软雅黑'
        # Text
        run2 = p.add_run()
        run2.text = item
        run2.font.size = Pt(font_size)
        run2.font.color.rgb = color
        run2.font.name = '微软雅黑'
        p.space_after = line_spacing
    return txBox

def add_two_column(slide, left_title, right_title, left_items, right_items,
                   top=Cm(5.5), col_width=Cm(13), font_size=15):
    """Add two-column layout"""
    # Left title
    add_text_box(slide, Cm(3.5), top, col_width, Cm(1), left_title,
                 font_size=18, color=GOLD, bold=True)
    add_gold_line(slide, Cm(3.5), top + Cm(0.8), Cm(5), Pt(2))
    add_bullet_list(slide, Cm(3.5), top + Cm(1.2), col_width, Cm(10),
                    left_items, font_size=font_size)
    # Right title
    add_text_box(slide, Cm(18.5), top, col_width, Cm(1), right_title,
                 font_size=18, color=GOLD, bold=True)
    add_gold_line(slide, Cm(18.5), top + Cm(0.8), Cm(5), Pt(2))
    add_bullet_list(slide, Cm(18.5), top + Cm(1.2), col_width, Cm(10),
                    right_items, font_size=font_size)

# ============================================================
# SLIDE 1: Title Slide
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)
# Gold accent bar at top
add_shape_bg(slide, Cm(0), Cm(0), Cm(33.867), Cm(0.5), GOLD)
# Title
add_text_box(slide, Cm(3.5), Cm(4.5), Cm(27), Cm(3),
             '佛教文化与宗教中国化研究', font_size=40,
             color=GOLD_LIGHT, bold=True, alignment=PP_ALIGN.CENTER)
# Gold line
add_gold_line(slide, Cm(10), Cm(8), Cm(14), Pt(3))
# Subtitle
add_text_box(slide, Cm(3.5), Cm(9), Cm(27), Cm(2),
             '历史进程 · 核心内涵 · 当代启示', font_size=24,
             color=LIGHT_GRAY, bold=False, alignment=PP_ALIGN.CENTER)
# Author & Date
add_text_box(slide, Cm(3.5), Cm(13), Cm(27), Cm(2),
             '汇报人：达真\n2026年6月15日', font_size=18,
             color=MID_GRAY, alignment=PP_ALIGN.CENTER)
# Bottom gold bar
add_shape_bg(slide, Cm(0), Cm(18.5), Cm(33.867), Cm(0.5), GOLD)

# ============================================================
# SLIDE 2: 导论 - 研究背景
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)
add_shape_bg(slide, Cm(0), Cm(0), Cm(33.867), Cm(0.5), GOLD)
add_text_box(slide, Cm(2), Cm(1.5), Cm(15), Cm(1.5),
             '导论：研究背景与意义', font_size=32,
             color=GOLD_LIGHT, bold=True)
add_gold_line(slide, Cm(2), Cm(3.2), Cm(6), Pt(3))

add_text_box(slide, Cm(2), Cm(4), Cm(30), Cm(1.5),
             '时代背景', font_size=22, color=GOLD, bold=True)
items_bg = [
    '党的二十大报告："坚持我国宗教中国化方向，积极引导宗教与社会主义社会相适应"',
    '全国宗教工作会议（2021）：将宗教中国化纳入文化自信自强总体布局',
    '佛教中国化是我国宗教中国化最成功的历史范例，跨越两千年的文化融合过程'
]
add_bullet_list(slide, Cm(2), Cm(5.5), Cm(30), Cm(4), items_bg, font_size=17)

add_text_box(slide, Cm(2), Cm(10), Cm(30), Cm(1),
             '研究意义', font_size=22, color=GOLD, bold=True)
items_sig = [
    '理论意义：构建"教义—制度—文化"三维分析框架，为宗教中国化提供理论模型',
    '实践意义：为当代各宗教中国化提供具体工作指南，尤其对道教理论现代化有借鉴价值',
    '跨学科视角：引入荣格分析心理学，探索心性论与现代心理学的对接路径'
]
add_bullet_list(slide, Cm(2), Cm(11), Cm(30), Cm(6), items_sig, font_size=17)

# ============================================================
# SLIDE 3: 研究方法
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)
add_shape_bg(slide, Cm(0), Cm(0), Cm(33.867), Cm(0.5), GOLD)
add_text_box(slide, Cm(2), Cm(1.5), Cm(15), Cm(1.5),
             '研究方法', font_size=32,
             color=GOLD_LIGHT, bold=True)
add_gold_line(slide, Cm(2), Cm(3.2), Cm(6), Pt(3))

methods = [
    ('历史文献分析法', '通过对佛教经典、历史文献、教理著作的梳理，还原佛教中国化的历史进程'),
    ('比较研究法', '通过佛教与道教中国化路径的比较，揭示宗教中国化的共性与个性规律'),
    ('跨学科研究法', '结合宗教学、哲学、心理学等多学科视角，引入荣格分析心理学理论框架'),
    ('案例分析法', '选取佛教中国化过程中的关键事件、代表人物和典型案例进行深度分析'),
]

for i, (title, desc) in enumerate(methods):
    y = Cm(4.5 + i * 3)
    # Number circle
    add_text_box(slide, Cm(3), y, Cm(2), Cm(1.5),
                 str(i+1), font_size=28, color=GOLD, bold=True)
    add_text_box(slide, Cm(5), y, Cm(12), Cm(1),
                 title, font_size=20, color=WHITE, bold=True)
    add_text_box(slide, Cm(5), y + Cm(1), Cm(25), Cm(1.5),
                 desc, font_size=16, color=TEXT_BODY)

# ============================================================
# SLIDE 4: 佛教传入的冲突与调适
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)
add_shape_bg(slide, Cm(0), Cm(0), Cm(33.867), Cm(0.5), GOLD)
add_text_box(slide, Cm(2), Cm(1.5), Cm(20), Cm(1.5),
             '一、佛教传入：冲突与调适的开端', font_size=30,
             color=GOLD_LIGHT, bold=True)
add_gold_line(slide, Cm(2), Cm(3.2), Cm(6), Pt(3))

add_text_box(slide, Cm(2), Cm(4.2), Cm(15), Cm(1),
             '核心文化冲突', font_size=22, color=GOLD, bold=True)

conflicts = [
    '出家 vs 孝道 — "身体发肤，受之父母，不敢毁伤"（《孝经》）',
    '沙门不敬王者 — 僧人不行君臣之礼，引发东晋政教论辩',
    '夷夏之辨 — 佛教被视为"夷狄之教"，南北朝时期激烈论争',
]
add_bullet_list(slide, Cm(2), Cm(5.5), Cm(15), Cm(8), conflicts, font_size=17)

add_text_box(slide, Cm(18.5), Cm(4.2), Cm(13), Cm(1),
             '初步调适策略', font_size=22, color=GOLD, bold=True)

adaptations = [
    '"格义"之法 — 以老庄概念比附佛理',
    '以"无"释"空"、以"无为"释"涅槃"',
    '依附黄老之学与玄学清谈',
    '译经活动奠定文献基础',
]
add_bullet_list(slide, Cm(18.5), Cm(5.5), Cm(13), Cm(8), adaptations, font_size=17)

# ============================================================
# SLIDE 5: 翻译中国化三阶段
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)
add_shape_bg(slide, Cm(0), Cm(0), Cm(33.867), Cm(0.5), GOLD)
add_text_box(slide, Cm(2), Cm(1.5), Cm(20), Cm(1.5),
             '二、翻译中国化：三阶段演进', font_size=30,
             color=GOLD_LIGHT, bold=True)
add_gold_line(slide, Cm(2), Cm(3.2), Cm(6), Pt(3))

# Create table
tbl_left = Cm(3)
tbl_top = Cm(5)
tbl_width = Cm(28)
tbl_height = Cm(10)
rows, cols = 4, 5
table_shape = slide.shapes.add_table(rows, cols, tbl_left, tbl_top, tbl_width, tbl_height)
table = table_shape.table

# Set column widths
table.columns[0].width = Cm(4)
table.columns[1].width = Cm(5)
table.columns[2].width = Cm(7)
table.columns[3].width = Cm(6)
table.columns[4].width = Cm(6)

headers = ['阶段', '代表人物', '翻译特点', '代表译著', '中国化意义']
data = [
    ['格义期\n(东汉)', '安世高\n支娄迦谶', '以老庄概念比附佛理', '《安般守意经》', '初步建立\n中印概念对应'],
    ['意译期\n(魏晋)', '鸠摩罗什', '不拘字面\n注重义理通达', '《金刚经》《法华经》', '语言与思想\n双重中国化'],
    ['新译期\n(唐代)', '玄奘', '精准严谨\n体系完整', '《大般若经》\n《成唯识论》', '系统化\n理论建构'],
]

for j, h in enumerate(headers):
    cell = table.cell(0, j)
    cell.text = h
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(16)
        p.font.color.rgb = BG_DARK
        p.font.bold = True
        p.font.name = '微软雅黑'
        p.alignment = PP_ALIGN.CENTER
    cell.fill.solid()
    cell.fill.fore_color.rgb = GOLD

for i, row_data in enumerate(data):
    for j, val in enumerate(row_data):
        cell = table.cell(i+1, j)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(15)
            p.font.color.rgb = TEXT_BODY
            p.font.name = '微软雅黑'
            p.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = BG_ACCENT if i % 2 == 0 else RGBColor(0x15, 0x25, 0x40)

# ============================================================
# SLIDE 6: 隋唐八大宗派
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)
add_shape_bg(slide, Cm(0), Cm(0), Cm(33.867), Cm(0.5), GOLD)
add_text_box(slide, Cm(2), Cm(1.5), Cm(20), Cm(1.5),
             '二、隋唐义理中国化：八大宗派', font_size=30,
             color=GOLD_LIGHT, bold=True)
add_gold_line(slide, Cm(2), Cm(3.2), Cm(6), Pt(3))

sects = [
    ('天台宗', '智顗 · "一念三千""三谛圆融"'),
    ('华严宗', '法藏 · "法界缘起""四法界"'),
    ('禅宗', '慧能 · "不立文字""直指人心""顿悟成佛"'),
    ('净土宗', '善导 · "称名念佛"，民间影响最大'),
    ('律宗', '道宣 · 系统整理佛教戒律'),
    ('法相宗', '玄奘 · 精密逻辑分析，过于学术化'),
    ('密宗', '善无畏 · 后融入其他宗派'),
    ('三论宗', '吉藏 · "八不中道"思想'),
]

for i, (name, desc) in enumerate(sects):
    col = i // 4
    row = i % 4
    x = Cm(3.5 + col * 14.5)
    y = Cm(4.5 + row * 2.8)
    # Name in gold
    add_text_box(slide, x, y, Cm(5), Cm(0.8), name,
                 font_size=20, color=GOLD, bold=True)
    # Description
    add_text_box(slide, x + Cm(5), y, Cm(9), Cm(1), desc,
                 font_size=16, color=TEXT_BODY)
    # Divider line
    add_gold_line(slide, x, y + Cm(1.2), Cm(13), Pt(1))

# ============================================================
# SLIDE 7: 心性论转化
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)
add_shape_bg(slide, Cm(0), Cm(0), Cm(33.867), Cm(0.5), GOLD)
add_text_box(slide, Cm(2), Cm(1.5), Cm(20), Cm(1.5),
             '核心成果：心性论的创造性转化', font_size=30,
             color=GOLD_LIGHT, bold=True)
add_gold_line(slide, Cm(2), Cm(3.2), Cm(6), Pt(3))

# Quote box
add_shape_bg(slide, Cm(3), Cm(4.5), Cm(28), Cm(3), BG_ACCENT)
add_text_box(slide, Cm(4), Cm(4.8), Cm(26), Cm(1),
             '"一切众生悉有佛性。"', font_size=24,
             color=GOLD_LIGHT, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Cm(4), Cm(6), Cm(26), Cm(1),
             '——《大般涅槃经》', font_size=16,
             color=MID_GRAY, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Cm(3), Cm(8.5), Cm(28), Cm(1),
             '从"缘起性空"到"明心见性"的理论转向', font_size=22,
             color=GOLD, bold=True)

turns = [
    '天台宗"性具"说 — 众生本性具足一切法',
    '华严宗"性起"说 — 一切法由心性而起',
    '禅宗"即心即佛" — 心就是佛，无需外求',
    '与儒家"性善论"、道家"自然本性论"形成思想共振',
]
add_bullet_list(slide, Cm(3), Cm(10), Cm(28), Cm(8), turns, font_size=18, line_spacing=Pt(10))

# ============================================================
# SLIDE 8: 制度与文化中国化
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)
add_shape_bg(slide, Cm(0), Cm(0), Cm(33.867), Cm(0.5), GOLD)
add_text_box(slide, Cm(2), Cm(1.5), Cm(20), Cm(1.5),
             '制度与文化的中国化', font_size=30,
             color=GOLD_LIGHT, bold=True)
add_gold_line(slide, Cm(2), Cm(3.2), Cm(6), Pt(3))

add_two_column(slide,
    '制度中国化', '文化中国化',
    [
        '《百丈清规》："一日不作，一日不食"',
        '农禅并重 — 打破印度乞食传统',
        '传法世系宗法化 — "一脉相承"',
        '僧官制度（唐）→ 度牒制度（宋）',
        '纳入国家管理体系',
    ],
    [
        '儒释道三教合一 — 宋明理学吸收佛教心性论',
        '契嵩《辅教篇》："孝为戒先"',
        '观音信仰本土化 — 从男菩萨到慈母形象',
        '深刻影响诗歌、绘画、书法、建筑、雕塑',
        '《西游记》《红楼梦》等文学巨著',
    ],
    top=Cm(4), font_size=16
)

# ============================================================
# SLIDE 9: 曲折与教训 - 三武一宗灭佛
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)
add_shape_bg(slide, Cm(0), Cm(0), Cm(33.867), Cm(0.5), GOLD)
add_text_box(slide, Cm(2), Cm(1.5), Cm(20), Cm(1.5),
             '三、曲折与教训：三武一宗灭佛', font_size=30,
             color=GOLD_LIGHT, bold=True)
add_gold_line(slide, Cm(2), Cm(3.2), Cm(6), Pt(3))

# Table
tbl_left = Cm(2.5)
tbl_top = Cm(4.5)
tbl_width = Cm(29)
tbl_height = Cm(7)
rows, cols = 5, 4
table_shape = slide.shapes.add_table(rows, cols, tbl_left, tbl_top, tbl_width, tbl_height)
table = table_shape.table

table.columns[0].width = Cm(6)
table.columns[1].width = Cm(7)
table.columns[2].width = Cm(8)
table.columns[3].width = Cm(8)

headers = ['事件', '时间/主导者', '直接原因', '深层原因']
data = [
    ['北魏太武帝灭佛', '446年 / 拓跋焘', '僧人参与叛乱', '势力膨胀威胁皇权'],
    ['北周武帝灭佛', '574年 / 宇文邕', '佛道论争', '儒道联合抵制 + 经济因素'],
    ['唐武宗灭佛\n(会昌法难)', '845年 / 李炎', '崇道贬佛', '寺院经济过度膨胀\n与国家争劳动力税收'],
    ['后周世宗灭佛', '955年 / 柴荣', '整顿寺院', '铜钱短缺 毁佛铸钱'],
]

for j, h in enumerate(headers):
    cell = table.cell(0, j)
    cell.text = h
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(15)
        p.font.color.rgb = BG_DARK
        p.font.bold = True
        p.font.name = '微软雅黑'
        p.alignment = PP_ALIGN.CENTER
    cell.fill.solid()
    cell.fill.fore_color.rgb = ACCENT_RED

for i, row_data in enumerate(data):
    for j, val in enumerate(row_data):
        cell = table.cell(i+1, j)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(14)
            p.font.color.rgb = TEXT_BODY
            p.font.name = '微软雅黑'
            p.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = BG_ACCENT if i % 2 == 0 else RGBColor(0x15, 0x25, 0x40)

# Key lesson
add_text_box(slide, Cm(2.5), Cm(12.5), Cm(29), Cm(2),
             '⚠ 深层共性：寺院经济过度膨胀 → 与国家利益冲突 → 宗教须妥善处理与世俗权力的关系',
             font_size=17, color=GOLD_LIGHT, bold=True)

# ============================================================
# SLIDE 10: 历史经验总结
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)
add_shape_bg(slide, Cm(0), Cm(0), Cm(33.867), Cm(0.5), GOLD)
add_text_box(slide, Cm(2), Cm(1.5), Cm(15), Cm(1.5),
             '五、历史经验与教训', font_size=32,
             color=GOLD_LIGHT, bold=True)
add_gold_line(slide, Cm(2), Cm(3.2), Cm(6), Pt(3))

experiences = [
    ('主动适应', '从"移植"到"扎根"——宗教中国化是主动的、内在的转化过程'),
    ('不变随缘', '保持核心信仰的同时实现形式创新——华严宗"不变随缘"的智慧'),
    ('服务社会', '"庄严国土、利乐有情"——与社会发展同频共振'),
    ('价值协调', '与主流价值观相协调——文化融合的最高境界'),
]

for i, (title, desc) in enumerate(experiences):
    y = Cm(4.5 + i * 3)
    add_text_box(slide, Cm(3.5), y, Cm(5), Cm(1), title,
                 font_size=22, color=GOLD, bold=True)
    add_text_box(slide, Cm(3.5), y + Cm(1), Cm(25), Cm(1.5),
                 desc, font_size=17, color=TEXT_BODY)

add_text_box(slide, Cm(3.5), Cm(16.5), Cm(27), Cm(2),
             '⚠ 教训：经济与社会责任的平衡 | 理论深度与大众接受的平衡 | 宗教自主与国家管理的平衡',
             font_size=15, color=MID_GRAY)

# ============================================================
# SLIDE 11: 当代实践
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)
add_shape_bg(slide, Cm(0), Cm(0), Cm(33.867), Cm(0.5), GOLD)
add_text_box(slide, Cm(2), Cm(1.5), Cm(20), Cm(1.5),
             '六、当代实践：从"人间佛教"到今日', font_size=30,
             color=GOLD_LIGHT, bold=True)
add_gold_line(slide, Cm(2), Cm(3.2), Cm(6), Pt(3))

# Three stages
stages = [
    ('第一阶段：太虚大师奠基\n(1920s-1940s)', [
        '"人成即佛成"核心命题',
        '佛教应关注现实人生，服务社会',
        '教义用现代语言重新阐释',
    ]),
    ('第二阶段：星云大师推广\n(1960s-2020s)', [
        '佛光山全球300+寺院道场',
        '创办大学、医院、美术馆',
        '"四给"理念：信心、欢喜、希望、方便',
    ]),
    ('第三阶段：当代大陆探索\n(2000s至今)', [
        '教义阐释中国化',
        '寺院管理规范化',
        '社会服务制度化',
    ]),
]

for i, (title, items) in enumerate(stages):
    x = Cm(2.5 + i * 10)
    y = Cm(4.5)
    add_shape_bg(slide, x - Cm(0.3), y - Cm(0.3), Cm(9), Cm(1.8), GOLD)
    add_text_box(slide, x, y, Cm(8.5), Cm(1.6), title,
                 font_size=16, color=BG_DARK, bold=True)
    add_bullet_list(slide, x, y + Cm(2), Cm(9), Cm(10), items,
                    font_size=16, line_spacing=Pt(10))

# ============================================================
# SLIDE 12: 佛道中国化比较
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)
add_shape_bg(slide, Cm(0), Cm(0), Cm(33.867), Cm(0.5), GOLD)
add_text_box(slide, Cm(2), Cm(1.5), Cm(20), Cm(1.5),
             '七、佛道中国化比较分析', font_size=30,
             color=GOLD_LIGHT, bold=True)
add_gold_line(slide, Cm(2), Cm(3.2), Cm(6), Pt(3))

# Table
tbl_left = Cm(3)
tbl_top = Cm(4.5)
tbl_width = Cm(28)
tbl_height = Cm(12)
rows, cols = 9, 3
table_shape = slide.shapes.add_table(rows, cols, tbl_left, tbl_top, tbl_width, tbl_height)
table = table_shape.table

table.columns[0].width = Cm(5)
table.columns[1].width = Cm(11)
table.columns[2].width = Cm(12)

headers = ['比较维度', '佛教中国化', '道教中国化（当代）']
data = [
    ['起点', '外来宗教的本土化适应', '本土宗教的现代化转型'],
    ['主要张力', '夷夏之辨、佛儒冲突', '传统与现代的冲突、科学与宗教'],
    ['核心转化', '从"缘起性空"到"明心见性"', '从"修炼成仙"到"身心和谐"'],
    ['制度创新', '禅林清规、宗法世系', '待完善'],
    ['社会功能', '从出世到入世（人间佛教）', '从个人修炼到社会服务'],
    ['理论体系化', '八大宗派，体系完备', '相对分散，缺乏系统化阐释'],
    ['文化融合', '儒释道三教合一', '待深入'],
    ['国际传播', '全球性宗教', '主要限于华人文化圈'],
]

for j, h in enumerate(headers):
    cell = table.cell(0, j)
    cell.text = h
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(15)
        p.font.color.rgb = BG_DARK
        p.font.bold = True
        p.font.name = '微软雅黑'
        p.alignment = PP_ALIGN.CENTER
    cell.fill.solid()
    cell.fill.fore_color.rgb = GOLD

for i, row_data in enumerate(data):
    for j, val in enumerate(row_data):
        cell = table.cell(i+1, j)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(14)
            p.font.color.rgb = TEXT_BODY
            p.font.name = '微软雅黑'
            p.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = BG_ACCENT if i % 2 == 0 else RGBColor(0x15, 0x25, 0x40)

# ============================================================
# SLIDE 13: 心性论借鉴 - 道教方法对照表
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)
add_shape_bg(slide, Cm(0), Cm(0), Cm(33.867), Cm(0.5), GOLD)
add_text_box(slide, Cm(2), Cm(1.5), Cm(22), Cm(1.5),
             '借鉴：道教修炼与现代心理学对照', font_size=30,
             color=GOLD_LIGHT, bold=True)
add_gold_line(slide, Cm(2), Cm(3.2), Cm(6), Pt(3))

# Table
tbl_left = Cm(2.5)
tbl_top = Cm(4.5)
tbl_width = Cm(29)
tbl_height = Cm(12)
rows, cols = 7, 4
table_shape = slide.shapes.add_table(rows, cols, tbl_left, tbl_top, tbl_width, tbl_height)
table = table_shape.table

table.columns[0].width = Cm(4)
table.columns[1].width = Cm(5)
table.columns[2].width = Cm(9)
table.columns[3].width = Cm(11)

headers = ['道教方法', '经典出处', '现代心理学对应', '应用场景']
data = [
    ['坐忘', '《庄子·大宗师》', '正念冥想、觉察训练', '焦虑缓解、压力管理'],
    ['心斋', '《庄子·人间世》', '冥想减压、情绪调节', '情绪障碍干预'],
    ['守一', '《太平经》', '注意力训练、专注力培养', 'ADHD辅助干预'],
    ['齐物', '《庄子·齐物论》', '认知重构、视角转换', '认知行为疗法'],
    ['虚静', '《道德经》第16章', '放松训练、身心整合', '失眠干预'],
    ['内观', '《黄庭经》', '身体扫描、躯体觉察', '躯体化症状干预'],
]

for j, h in enumerate(headers):
    cell = table.cell(0, j)
    cell.text = h
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(15)
        p.font.color.rgb = BG_DARK
        p.font.bold = True
        p.font.name = '微软雅黑'
        p.alignment = PP_ALIGN.CENTER
    cell.fill.solid()
    cell.fill.fore_color.rgb = ACCENT_GREEN

for i, row_data in enumerate(data):
    for j, val in enumerate(row_data):
        cell = table.cell(i+1, j)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(14)
            p.font.color.rgb = TEXT_BODY
            p.font.name = '微软雅黑'
            p.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = BG_ACCENT if i % 2 == 0 else RGBColor(0x15, 0x25, 0x40)

# ============================================================
# SLIDE 14: 道教现代化路径
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)
add_shape_bg(slide, Cm(0), Cm(0), Cm(33.867), Cm(0.5), GOLD)
add_text_box(slide, Cm(2), Cm(1.5), Cm(20), Cm(1.5),
             '道教理论现代化：四条路径', font_size=30,
             color=GOLD_LIGHT, bold=True)
add_gold_line(slide, Cm(2), Cm(3.2), Cm(6), Pt(3))

paths = [
    ('经典新释', '对《道德经》《庄子》作出符合当代社会需求的新阐释\n如："致虚极，守静笃" → 减压与情绪管理方法'),
    ('概念转换', '将核心概念（道、德、气、阴阳、性命）转化为当代学术话语\n如："气" → "生命能量"或"身心状态"'),
    ('实践创新', '将修炼方法（静坐、导引、斋醮）转化为现代人可接受的身心实践\n建立道教心理疏导服务体系'),
    ('价值对接', '将道教核心价值观与社会主义核心价值观对接\n如："道法自然" ↔ "生态文明建设"'),
]

for i, (title, desc) in enumerate(paths):
    x = Cm(3) if i < 2 else Cm(18)
    y = Cm(4.5) if i < 2 else Cm(11)
    # Box background
    add_shape_bg(slide, x - Cm(0.3), y - Cm(0.3), Cm(13.5), Cm(4.5), BG_ACCENT)
    add_gold_line(slide, x, y, Cm(3), Pt(2))
    add_text_box(slide, x + Cm(3.5), y - Cm(0.1), Cm(9), Cm(0.8), title,
                 font_size=20, color=GOLD, bold=True)
    add_text_box(slide, x, y + Cm(0.8), Cm(13), Cm(3.5), desc,
                 font_size=15, color=TEXT_BODY, line_spacing=1.4)

# ============================================================
# SLIDE 15: 结论
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)
add_shape_bg(slide, Cm(0), Cm(0), Cm(33.867), Cm(0.5), GOLD)
add_text_box(slide, Cm(2), Cm(1.5), Cm(15), Cm(1.5),
             '结论', font_size=32,
             color=GOLD_LIGHT, bold=True)
add_gold_line(slide, Cm(2), Cm(3.2), Cm(6), Pt(3))

conclusions = [
    '佛教中国化是以教义思想为核心、制度仪轨为保障、文化艺术为载体、服务社会为归宿的完整过程',
    '历史教训提醒我们处理好三组关系：经济独立与社会责任 | 理论深度与大众接受 | 宗教自主与国家管理',
    '道教可借鉴佛教经验：心性论对接现代心理学 | 社会功能拓展为心理疏导 | 文化传播利用新技术',
    '宗教中国化是宗教自身发展的内在规律，只有自觉顺应，才能在当代社会发挥积极作用',
]
add_bullet_list(slide, Cm(2.5), Cm(4.5), Cm(29), Cm(10), conclusions, font_size=19, line_spacing=Pt(16))

# Final quote
add_shape_bg(slide, Cm(5), Cm(14), Cm(24), Cm(3), BG_ACCENT)
add_text_box(slide, Cm(6), Cm(14.5), Cm(22), Cm(2),
             '"道教作为中国本土宗教，推动教义理论的现代化阐释\n和社会服务功能的拓展，是实现当代转型的必由之路。"',
             font_size=18, color=GOLD_LIGHT, bold=True, alignment=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 16: 致谢
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)
add_shape_bg(slide, Cm(0), Cm(0), Cm(33.867), Cm(0.5), GOLD)
add_text_box(slide, Cm(3.5), Cm(5), Cm(27), Cm(3),
             '感谢聆听', font_size=44,
             color=GOLD_LIGHT, bold=True, alignment=PP_ALIGN.CENTER)
add_gold_line(slide, Cm(10), Cm(9), Cm(14), Pt(3))
add_text_box(slide, Cm(3.5), Cm(10.5), Cm(27), Cm(2),
             '佛教文化与宗教中国化研究\n汇报人：达真 · 2026年6月15日', font_size=20,
             color=MID_GRAY, alignment=PP_ALIGN.CENTER)
add_shape_bg(slide, Cm(0), Cm(18.5), Cm(33.867), Cm(0.5), GOLD)

# ============================================================
# Save
# ============================================================
output_path = '/home/admin/.openclaw/workspace/佛教文化与宗教中国化报告.pptx'
prs.save(output_path)
print(f'Saved: {output_path}')
