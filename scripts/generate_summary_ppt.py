#!/usr/bin/env python3
"""小龙虾网络框架运行机制与模块学习总结 - 汇报 PPT 生成器"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── 配色方案 ──
COLOR_DARK = RGBColor(0x1A, 0x23, 0x3A)       # 深蓝黑
COLOR_PRIMARY = RGBColor(0xE8, 0x4D, 0x2A)    # 龙虾红
COLOR_SECONDARY = RGBColor(0x2C, 0x5F, 0x8A)  # 稳重蓝
COLOR_ACCENT = RGBColor(0xF5, 0xA6, 0x23)     # 暖橙
COLOR_GREEN = RGBColor(0x27, 0xAE, 0x60)       # 成功绿
COLOR_RED = RGBColor(0xC0, 0x39, 0x2B)         # 警告红
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_LIGHT_BG = RGBColor(0xF5, 0xF7, 0xFA)
COLOR_GRAY = RGBColor(0x7F, 0x8C, 0x8D)
COLOR_TEXT = RGBColor(0x2C, 0x3E, 0x50)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── 工具函数 ──

def add_bg(slide, color=COLOR_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, color, alpha=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if alpha is not None:
        shape.fill.fore_color.brightness = alpha
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18, color=COLOR_TEXT, bold=False, alignment=PP_ALIGN.LEFT, font_name='Microsoft YaHei'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_title_bar(slide, title_text, subtitle_text=None):
    """标准标题栏"""
    add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), COLOR_DARK)
    # 龙虾红装饰线
    add_rect(slide, Inches(0), Inches(1.2), Inches(13.333), Inches(0.06), COLOR_PRIMARY)
    add_text_box(slide, Inches(0.8), Inches(0.15), Inches(11), Inches(0.7), title_text, 26, COLOR_WHITE, True)
    if subtitle_text:
        add_text_box(slide, Inches(0.8), Inches(0.75), Inches(11), Inches(0.4), subtitle_text, 14, COLOR_GRAY)

def add_table(slide, left, top, width, height, rows, cols, data, col_widths=None, header_color=COLOR_SECONDARY):
    """添加格式化表格"""
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = str(data[r][c]) if data[r][c] is not None else ''
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(13)
                paragraph.font.name = 'Microsoft YaHei'
                paragraph.alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

            if r == 0:  # Header
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.color.rgb = COLOR_WHITE
                    paragraph.font.bold = True
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLOR_LIGHT_BG if r % 2 == 1 else COLOR_WHITE
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.color.rgb = COLOR_TEXT

    return table_shape

def add_bullet_points(slide, left, top, width, height, items, font_size=15, color=COLOR_TEXT, bullet_char='▸'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f'{bullet_char} {item}'
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = 'Microsoft YaHei'
        p.space_after = Pt(6)
    return txBox

def add_status_badge(slide, left, top, text, status='active'):
    """状态标签"""
    colors = {'active': COLOR_GREEN, 'warning': COLOR_ACCENT, 'error': COLOR_RED, 'info': COLOR_SECONDARY, 'pending': COLOR_GRAY}
    w = Inches(1.0)
    h = Inches(0.35)
    shape = add_rect(slide, left, top, w, h, colors.get(status, COLOR_GRAY))
    shape.text_frame.word_wrap = True
    p = shape.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(10)
    p.font.color.rgb = COLOR_WHITE
    p.font.bold = True
    p.font.name = 'Microsoft YaHei'
    p.alignment = PP_ALIGN.CENTER

# ════════════════════════════════════════════════════════════
# Slide 1: 封面
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, COLOR_DARK)
# 龙虾红装饰带
add_rect(slide, Inches(0), Inches(2.5), Inches(13.333), Inches(0.06), COLOR_PRIMARY)
# 标题
add_text_box(slide, Inches(1.5), Inches(1.2), Inches(10), Inches(1.3),
             '小龙虾网络', 48, COLOR_WHITE, True, PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.5), Inches(2.8), Inches(10), Inches(1.0),
             '框架运行机制与模块学习总结', 32, COLOR_WHITE, True, PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.5), Inches(4.0), Inches(10), Inches(0.6),
             'V4.0+ 整体复盘汇报', 22, COLOR_PRIMARY, False, PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.5), Inches(4.8), Inches(10), Inches(0.5),
             '2026年7月10日  |  WorkBuddy 编制', 16, COLOR_GRAY, False, PP_ALIGN.CENTER)
# 底部装饰
add_rect(slide, Inches(0), Inches(7.1), Inches(13.333), Inches(0.4), COLOR_PRIMARY)
add_text_box(slide, Inches(1.5), Inches(7.15), Inches(10), Inches(0.3),
             '🦞 你不停对话，世界就不停扩展', 14, COLOR_WHITE, False, PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════
# Slide 2: 项目概述
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, COLOR_WHITE)
add_title_bar(slide, '项目概述', '小龙虾网络 Lobster Network')

add_text_box(slide, Inches(0.8), Inches(1.6), Inches(11), Inches(0.7),
             '基于「对话即创造」理论的多 Agent 协作网络 — 将哲学命题工程化为可运行系统',
             17, COLOR_GRAY, False)

overview_data = [
    ['指标', '数值', '指标', '数值'],
    ['GitHub 仓库', 'zhugebin-hub/lobster-network', '节点数', '6 (1coach+4student+1observer)'],
    ['当前版本', 'V4.0+ (主力) / v0.4.1 (部署)', '学习引擎', '14 个领域引擎'],
    ['CC 消息追踪', '13 completed / 1 escalated', '7月提交数', '202+'],
    ['心跳保活', '每小时自动执行', '成功率', '100%'],
]
add_table(slide, Inches(0.8), Inches(2.5), Inches(11.5), Inches(3.5), 5, 4, overview_data,
          col_widths=[Inches(2.2), Inches(3.5), Inches(2.2), Inches(3.6)])

add_text_box(slide, Inches(0.8), Inches(6.3), Inches(11), Inches(0.5),
             '核心理念：一人一世界观 × 世界是对话 × 世界是编程的', 15, COLOR_PRIMARY, True, PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════
# Slide 3: 架构全景
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, COLOR_WHITE)
add_title_bar(slide, '架构全景图', '五层分层架构：应用 → 运营 → 同步通信 → 框架 → 仪表盘')

layers = [
    ('应用层 Domains', '围棋 · 海报 · 学习矩阵 · 药物发现 · 炒股 · 论文写作', COLOR_PRIMARY),
    ('运营层 Core', '任务调度 · 学生 Agent · 教练系统 · 监控工具', COLOR_ACCENT),
    ('同步通信层', 'sync_manager V4.0 · CC Protocol · 消息队列 · SSH 通道 V2', COLOR_SECONDARY),
    ('框架层 Framework', 'Node · Dialogue · Emergence · WorldState · IndraNet · TimeArbitrage', COLOR_GREEN),
    ('仪表盘层 Dashboard V5.0', '网络拓扑 · 系统监控 · 学习进度 · 项目指挥中心', RGBColor(0x8E, 0x44, 0xAD)),
]
y = Inches(1.6)
for name, desc, color in layers:
    add_rect(slide, Inches(1.5), y, Inches(10.3), Inches(0.95), color)
    add_text_box(slide, Inches(1.7), y + Inches(0.1), Inches(3.5), Inches(0.35), name, 16, COLOR_WHITE, True)
    add_text_box(slide, Inches(1.7), y + Inches(0.5), Inches(9.8), Inches(0.35), desc, 12, COLOR_WHITE, False)
    y += Inches(1.08)

# ════════════════════════════════════════════════════════════
# Slide 4: 节点拓扑
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, COLOR_WHITE)
add_title_bar(slide, '节点拓扑与注册表', '6 节点协同：1 coach + 4 student + 1 observer')

node_data = [
    ['节点 ID', '名称', '角色', '服务器', '版本', '状态'],
    ['zhugema', '诸葛马', 'coach (主节点)', '47.93.6.57', '0.5.0', 'active ✅'],
    ['xiaochen', '小陈', 'student (稳健型)', '121.43.80.231', '0.5.0', 'active ✅'],
    ['zhuguxia', '诸葛虾', 'student (加速型)', '60.205.139.51', '0.6.0', 'active ✅'],
    ['qoder', '小龙虾', 'student (实战型)', '本地 Mac', '0.6.0', 'active ✅'],
    ['xiaowei', '小薇', 'student (实战型)', 'local', '0.6.0', 'active ✅'],
    ['zhugebin-001', '诸葛斌', 'observer (发起人)', '—', '0.5.0', 'active ✅'],
]
add_table(slide, Inches(0.5), Inches(1.6), Inches(12.3), Inches(3.8), 7, 6, node_data,
          col_widths=[Inches(2), Inches(1.6), Inches(2.5), Inches(2.4), Inches(1.4), Inches(2)])

# 双注册表说明
add_text_box(slide, Inches(0.5), Inches(5.7), Inches(12), Inches(0.4),
             '双注册表机制', 16, COLOR_SECONDARY, True)

items = [
    '服务器注册表 /shared/registry/registry.json @ 121.43.80.231 — NFS 共享，通过 SSH 访问',
    '本地注册表 .shared/messages/ @ Claw/lobster-network — sync_manager V4.0 管理，含全量 6 节点',
    '当前服务器端在线节点：qoder (active) + workbuddy (active)',
]
add_bullet_points(slide, Inches(0.5), Inches(6.1), Inches(12), Inches(1.2), items, 13, COLOR_TEXT)

# ════════════════════════════════════════════════════════════
# Slide 5: 通信协议
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, COLOR_WHITE)
add_title_bar(slide, '消息通信与同步机制', '两层消息系统 + CC Protocol + 心跳保活')

# L1
add_rect(slide, Inches(0.5), Inches(1.5), Inches(5.8), Inches(2.7), COLOR_LIGHT_BG)
add_text_box(slide, Inches(0.7), Inches(1.55), Inches(5.4), Inches(0.4), 'L1 简单消息 - 心跳格式', 15, COLOR_SECONDARY, True)
add_bullet_points(slide, Inches(0.7), Inches(2.0), Inches(5.4), Inches(2.0), [
    '路径: /shared/messages/from-<node>/',
    '用途: 节点保活、状态广播',
    '频率: 每小时 1 次 (WorkBuddy 自动化)',
    '执行方式: SSH stdin 管道',
    '保留最近 5 个心跳文件',
    '7月7日至今成功率 100%',
], 12, COLOR_TEXT)

# L2
add_rect(slide, Inches(6.8), Inches(1.5), Inches(5.8), Inches(2.7), COLOR_LIGHT_BG)
add_text_box(slide, Inches(7.0), Inches(1.55), Inches(5.4), Inches(0.4), 'L2 可靠消息 - CC Protocol v1.1', 15, COLOR_SECONDARY, True)
add_bullet_points(slide, Inches(7.0), Inches(2.0), Inches(5.4), Inches(2.0), [
    '路径: .shared/messages/queue/<node>/inbox/',
    '生命周期: 发送 → Ack → cc_tracking (completed)',
    '超时处理: cc_escalate_expired.py → escalated',
    '当前追踪: 13 completed / 1 escalated',
    '支持重试和故障切换',
], 12, COLOR_TEXT)

# sync_manager
add_rect(slide, Inches(0.5), Inches(4.5), Inches(12.1), Inches(1.8), RGBColor(0xE8, 0xF0, 0xFE))
add_text_box(slide, Inches(0.7), Inches(4.55), Inches(11.5), Inches(0.4), 'sync_manager V4.0 核心功能', 15, COLOR_SECONDARY, True)
add_bullet_points(slide, Inches(0.7), Inches(5.0), Inches(5.5), Inches(1.2), [
    '统一路径验证 — 标准化目录结构',
    '消息格式校验 — JSON 合规性检查',
], 12, COLOR_TEXT)
add_bullet_points(slide, Inches(7.0), Inches(5.0), Inches(5.5), Inches(1.2), [
    '同步状态追踪 — 实时记录各节点状态',
    '自动修复 — 检测并修复常见同步问题',
], 12, COLOR_TEXT)

# ════════════════════════════════════════════════════════════
# Slide 6: 学习引擎矩阵
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, COLOR_WHITE)
add_title_bar(slide, '学习引擎矩阵 (14 个领域引擎)', '已完成 3 个 · 进行中 4 个 · 规划中 7 个')

lm_data = [
    ['#', '模块', '状态', '题数', '训练器'],
    ['1', '世界杯预测', '✅ 完成', '~30', '✅'],
    ['2', '炒股预测 V1.0', '✅ 完成', '120', '✅'],
    ['3', '网络协议', '✅ 完成', '90', '✅'],
    ['4', '药物发现', '🔄 Day1 完成', '—', '🔄'],
    ['5', '论文写作', '🔄 进行中', '—', '🔄'],
    ['6', 'AI/ML · 网络安全 · 数据结构 · 围棋 · 网络工程', '📋 规划中', '—', '📋'],
    ['7', '信号竞技场 · 海报 · 通用逻辑', '📋 规划中', '—', '📋'],
]
add_table(slide, Inches(0.5), Inches(1.5), Inches(12.3), Inches(3.8), 8, 5, lm_data,
          col_widths=[Inches(0.6), Inches(3.5), Inches(2.2), Inches(1.8), Inches(1.8)])

# 四层反馈循环
add_text_box(slide, Inches(0.5), Inches(5.5), Inches(5.8), Inches(0.4), '四层反馈循环', 15, COLOR_SECONDARY, True)
feedback_items = [
    'L1 即时反馈：每个任务完成后自动评估',
    'L2 每日反馈：教练分析表现，调整次日计划',
    'L3 每周反馈：跨学生讨论赛、复盘学习',
    'L4 任务反馈：人类灵魂提供方向和验收',
]
add_bullet_points(slide, Inches(0.5), Inches(5.9), Inches(5.8), Inches(1.5), feedback_items, 12, COLOR_TEXT)

add_text_box(slide, Inches(7.0), Inches(5.5), Inches(5.8), Inches(0.4), '五维时间套利引擎', 15, COLOR_SECONDARY, True)
time_items = [
    '速率套利：快节点生成 → 慢节点深化',
    '错峰套利：00:00-06:00 低成本算力',
    '反思套利：艾宾浩斯遗忘曲线（3日复习）',
    '复利套利：E_total = E₁ × (1+r)^(N-1)',
    '时距套利：48-72h 知识价值峰值',
]
add_bullet_points(slide, Inches(7.0), Inches(5.9), Inches(5.8), Inches(1.5), time_items, 12, COLOR_TEXT)

# ════════════════════════════════════════════════════════════
# Slide 7: 学员训练状态
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, COLOR_WHITE)
add_title_bar(slide, '学员训练状态', '当前训练活跃度需关注 — 多个模块连续 12 天无新训练')

st_data = [
    ['学员', '模块', '当前阶段', '最后训练日期', '状态'],
    ['xiaochen (小陈)', '网络协议', 'phase1', '—', '🔄 进行中'],
    ['xiaochen (小陈)', '炒股预测', '—', '2026-06-26', '⏸️ 停滞 14天'],
    ['zhuguxia (诸葛虾)', '炒股预测', '—', '2026-06-26', '⏸️ 停滞 14天'],
    ['zhugebin-001 (诸葛斌)', '炒股预测', '—', '2026-06-26', '⏸️ 停滞 14天'],
    ['zhugebin-001 (诸葛斌)', '论文写作', '—', '2026-07-05', '🔄 最近活跃'],
    ['qoder', '网络协议', 'phase1', '—', '🔄 进行中'],
]
add_table(slide, Inches(0.5), Inches(1.5), Inches(12.3), Inches(3.6), 7, 5, st_data,
          col_widths=[Inches(3), Inches(2.5), Inches(2), Inches(2.5), Inches(2.3)])

# Highlight
add_rect(slide, Inches(0.5), Inches(5.4), Inches(12.1), Inches(1.8), RGBColor(0xFF, 0xF3, 0xCD))
add_text_box(slide, Inches(0.7), Inches(5.45), Inches(11.5), Inches(0.35), '⚠️ 关键发现', 15, COLOR_ACCENT, True)
findings = [
    '炒股预测模块 3 个学员（小陈/诸葛虾/诸葛斌）全部停滞，最后训练于 2026-06-26，距今已 14 天',
    '仅网络协议和论文写作模块有最近活跃，其他引擎尚未配备完整训练器',
    '建议优先恢复 P1 方案：7 天内重启炒股预测 + 论文写作训练',
]
add_bullet_points(slide, Inches(0.7), Inches(5.8), Inches(11.5), Inches(1.3), findings, 13, COLOR_RED)

# ════════════════════════════════════════════════════════════
# Slide 8: 药物发现项目
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, COLOR_WHITE)
add_title_bar(slide, '药物发现项目 · Day1 里程碑', '2026-07-09 启动 | 食物过敏防治药物研制科学智能体')

# Day1 成果
add_rect(slide, Inches(0.5), Inches(1.5), Inches(5.8), Inches(2.2), RGBColor(0xE8, 0xF8, 0xF5))
add_text_box(slide, Inches(0.7), Inches(1.55), Inches(5.4), Inches(0.4), 'Day1 研究成果', 15, COLOR_GREEN, True)
day1_items = [
    '6 个过敏原已识别',
    '6 个药物靶点已确认',
    '8 个候选化合物已筛选',
    '6 个临床设计方案已生成',
]
add_bullet_points(slide, Inches(0.7), Inches(2.0), Inches(5.4), Inches(1.5), day1_items, 13, COLOR_GREEN, '✓')

# 研究目标
add_rect(slide, Inches(6.8), Inches(1.5), Inches(5.8), Inches(2.2), RGBColor(0xF0, 0xF4, 0xFA))
add_text_box(slide, Inches(7.0), Inches(1.55), Inches(5.4), Inches(0.4), '研究目标 (Phase 1)', 15, COLOR_SECONDARY, True)
goal_items = [
    '构建食物过敏原知识图谱 (10万+节点)',
    '发现并评估药物靶点 (10个关键靶点) ✅已识别6个',
    '搭建虚拟筛选管线 (百万级化合物)',
    '候选化合物优先级排序 (Top 20) ✅已识别8个',
    '临床前评估 (药代动力学+毒性) ⏳待启动',
]
add_bullet_points(slide, Inches(7.0), Inches(2.0), Inches(5.4), Inches(1.5), goal_items, 11.5, COLOR_TEXT)

# 节点分工
div_data = [
    ['节点', '角色', '研究方向'],
    ['诸葛马', '总教练/评审', '药物安全 + 监管 + 高级评审'],
    ['qoder', '计算化学专家', '知识图谱 + 虚拟筛选'],
    ['小陈', '免疫学专家', '过敏机制 + 靶点调研'],
    ['诸葛虾', '工具链专家', '可视化 + 工具开发'],
    ['诸葛斌', '发起人', '全流程 + 临床试验设计'],
    ['小薇', '实战型', '免疫疗法 + 临床执行'],
]
add_table(slide, Inches(0.5), Inches(4.0), Inches(12.3), Inches(3.2), 7, 3, div_data,
          col_widths=[Inches(2.5), Inches(3.5), Inches(6.3)])

# ════════════════════════════════════════════════════════════
# Slide 9: 今日运行复盘
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, COLOR_WHITE)
add_title_bar(slide, '今日运行复盘 (2026-07-09 ~ 07-10)', '24 小时运行指标 + 心跳保活 + 服务器注册表')

# KPI 卡片
kpi_data = [
    ('心跳执行', '26 次', '成功率 100%', COLOR_GREEN),
    ('Git 提交', '15 次', '高度活跃', COLOR_SECONDARY),
    ('CC 消息', '0 pending', '13 completed', COLOR_GREEN),
    ('训练提交', '连续12天', '无新训练 ❌', COLOR_RED),
    ('药物发现', 'Day1 完成', '里程碑 ✅', COLOR_ACCENT),
]
for i, (label, value, sub, color) in enumerate(kpi_data):
    x = Inches(0.5 + i * 2.5)
    add_rect(slide, x, Inches(1.5), Inches(2.3), Inches(1.7), color)
    add_text_box(slide, x + Inches(0.1), Inches(1.55), Inches(2.1), Inches(0.4), label, 14, COLOR_WHITE, True, PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.1), Inches(1.95), Inches(2.1), Inches(0.5), value, 22, COLOR_WHITE, True, PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.1), Inches(2.5), Inches(2.1), Inches(0.4), sub, 11, COLOR_WHITE, False, PP_ALIGN.CENTER)

# 心跳记录
add_text_box(slide, Inches(0.5), Inches(3.5), Inches(5.8), Inches(0.4), '最近 5 次心跳', 15, COLOR_SECONDARY, True)
hb_items = [
    '[07-10 01:29] HEARTBEAT workbuddy active | 保留5个心跳',
    '[07-10 00:34] HEARTBEAT workbuddy active | 保留5个心跳',
    '[07-09 23:38] HEARTBEAT workbuddy active | 保留5个心跳',
    '[07-09 22:42] HEARTBEAT workbuddy active | 保留5个心跳',
    '[07-09 21:47] HEARTBEAT workbuddy active | 保留5个心跳',
]
add_bullet_points(slide, Inches(0.5), Inches(3.9), Inches(5.8), Inches(2.0), hb_items, 11, COLOR_TEXT)

# 服务器状态
add_text_box(slide, Inches(7.0), Inches(3.5), Inches(5.8), Inches(0.4), '服务器注册表实时状态', 15, COLOR_SECONDARY, True)
reg_items = [
    'qoder → active | 最后心跳: 2026-07-10T01:42:10',
    'workbuddy → active | 最后心跳: 2026-07-10T01:29:52',
    '自动化: automation-1782283296487 持续运行中',
    '7月7日至今执行 50+ 次，全部成功',
]
add_bullet_points(slide, Inches(7.0), Inches(3.9), Inches(5.8), Inches(2.0), reg_items, 11, COLOR_TEXT)

# Git 活动
add_text_box(slide, Inches(0.5), Inches(5.8), Inches(12.1), Inches(0.4), 'Git 提交活动 (07-09~10 共 15 次)', 15, COLOR_SECONDARY, True)
git_items = [
    'docs: 小龙虾网络框架运行机制与模块学习总结 · 药物发现1.0 · 仪表盘V5.0 · 心跳部署 · 双平台同步',
    '三端 HEAD 一致 (本地 / GitHub origin / 服务器 server): a1e2287f',
]
add_bullet_points(slide, Inches(0.5), Inches(6.2), Inches(12), Inches(1.0), git_items, 12, COLOR_TEXT)

# ════════════════════════════════════════════════════════════
# Slide 10: 问题诊断
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, COLOR_WHITE)
add_title_bar(slide, '问题诊断 (3 个高优 + 5 个低优)', '涵盖训练活跃度、消息积压、版本碎片化、仪表盘滞后等')

issue_data = [
    ['#', '问题', '严重度', '影响'],
    ['P1', '连续 12 天无新训练提交', '🔴 高', '学习模块活跃度归零'],
    ['P2', '消息积压：诸葛马 inbox 27 条、qoder 29 条', '🟡 中', '消息处理滞后'],
    ['P3', 'CC 追踪有 1 条 escalated 消息未处理', '🟡 中', '消息可靠性降级'],
    ['P4', 'v0.4.1 部署版与 V4.0+ 主力版差异大', '🟡 中', '版本碎片化'],
    ['P5', '仪表盘快照滞后（最后扫描 07-08）', '🟢 低', '监控数据不实时'],
    ['P6', '学习状态文件字段不统一', '🟢 低', '状态解析困难'],
    ['P7', 'setup.py 版本号仍为 0.3.0', '🟢 低', '打包信息不准确'],
    ['P8', 'requirements.txt 缺 paramiko 依赖', '🟢 低', '安装不完整'],
]
add_table(slide, Inches(0.5), Inches(1.5), Inches(12.3), Inches(4.8), 9, 4, issue_data,
          col_widths=[Inches(0.6), Inches(4.5), Inches(1.8), Inches(5.4)])

add_text_box(slide, Inches(0.5), Inches(6.5), Inches(12), Inches(0.5),
             '⚠️ P1 为最高优先级：连续12天无新训练，3个学员停滞。需立即恢复炒股预测和论文写作训练。', 14, COLOR_RED, True)

# ════════════════════════════════════════════════════════════
# Slide 11: 优化方案与排期
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, COLOR_WHITE)
add_title_bar(slide, '优化方案与排期', '7 项方案，从 P0 立即执行到 P3 5 天内完成')

plan_data = [
    ['优先级', '方案', '目标', '预期完成', '负责人'],
    ['P0', '消息积压清理', '各节点 inbox ≤ 5 条', '立即', 'WorkBuddy'],
    ['P0', 'CC 升级消息处理', 'escalated 清零', '立即', 'WorkBuddy'],
    ['P1 🔴', '恢复学习训练', '3 个模块恢复活跃', '3 天内', '各学员节点'],
    ['P1', '打包修复', 'setup.py + requirements 同步', '1 天内', 'WorkBuddy'],
    ['P2', '版本统一', '部署版 → main 分支', '2 天内', 'WorkBuddy'],
    ['P2', '仪表盘实时化', '每 6 小时快照', '3 天内', 'WorkBuddy'],
    ['P3', '学习状态标准化', '统一 schema', '5 天内', 'WorkBuddy'],
]
add_table(slide, Inches(0.5), Inches(1.5), Inches(12.3), Inches(4.2), 8, 5, plan_data,
          col_widths=[Inches(1.2), Inches(3.2), Inches(3.5), Inches(2), Inches(2.4)])

add_text_box(slide, Inches(0.5), Inches(6.0), Inches(12), Inches(0.4), '恢复学习训练 (P1) 详细措施', 15, COLOR_SECONDARY, True)
measures = [
    '检查各学员训练器运行状态（SSH 到各服务器验证）',
    '重置停滞学员的训练状态，从断点继续',
    '启动每日训练调度器 automated_daily_training.py 定时自动化',
    '优先恢复炒股预测 (120题) 和论文写作模块 (已有训练器实现)',
]
add_bullet_points(slide, Inches(0.5), Inches(6.3), Inches(12), Inches(1.0), measures, 12, COLOR_TEXT)

# ════════════════════════════════════════════════════════════
# Slide 12: 路线图
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, COLOR_WHITE)
add_title_bar(slide, '版本路线图', 'V4.1 → V4.2 → V4.3 → V5.0 → V5.1')

roadmap = [
    ('V4.1', '打包修复 + 版本统一\n+ 消息清理', '2026-07-12', COLOR_GRAY),
    ('V4.2', '学习训练恢复\n+ 仪表盘实时化', '2026-07-15', COLOR_SECONDARY),
    ('V4.3', '学习状态标准化\n+ 监控增强', '2026-07-20', COLOR_SECONDARY),
    ('V5.0', '药物发现项目\nPhase 2 完成', '2026-07-30', COLOR_ACCENT),
    ('V5.1', '全量学习模块上线\n(14 个引擎)', '2026-08-15', COLOR_PRIMARY),
]
for i, (ver, desc, date, color) in enumerate(roadmap):
    x = Inches(0.8 + i * 2.5)
    # 版本圆圈
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.8), Inches(1.8), Inches(0.9), Inches(0.9))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    circle.text_frame.word_wrap = True
    p = circle.text_frame.paragraphs[0]
    p.text = ver
    p.font.size = Pt(18)
    p.font.color.rgb = COLOR_WHITE
    p.font.bold = True
    p.font.name = 'Microsoft YaHei'
    p.alignment = PP_ALIGN.CENTER

    # 描述
    add_text_box(slide, x, Inches(2.9), Inches(2.5), Inches(1.0), desc, 13, COLOR_TEXT, False, PP_ALIGN.CENTER)
    # 日期
    add_text_box(slide, x, Inches(3.8), Inches(2.5), Inches(0.4), date, 11, color, True, PP_ALIGN.CENTER)

    # 连接箭头
    if i < len(roadmap) - 1:
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x + Inches(2.5), Inches(2.1), Inches(0.35), Inches(0.3))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = COLOR_GRAY
        arrow.line.fill.background()

# 关键文件索引
add_text_box(slide, Inches(0.5), Inches(4.5), Inches(12), Inches(0.4), '关键文件索引', 15, COLOR_SECONDARY, True)
file_items = [
    '主力代码库: /Users/zgb/WorkBuddy/Claw/lobster-network (V4.0+ main分支)',
    '同步管理器: .shared/messages/sync_manager.py | CC追踪: .shared/messages/cc_tracking.json',
    '仪表盘: dashboard/status_snapshot.json | 心跳: scripts/heartbeat_workbuddy.py',
    '学习引擎: domains/learning/problems/*_engine.py | 药物发现: domains/drug-discovery/',
    '服务器注册表: /shared/registry/registry.json @ 121.43.80.231 (NFS共享)',
]
add_bullet_points(slide, Inches(0.5), Inches(4.9), Inches(12), Inches(2.5), file_items, 12, COLOR_TEXT)

# ════════════════════════════════════════════════════════════
# Slide 13: 总结与下一步
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, COLOR_DARK)
add_rect(slide, Inches(0), Inches(3.0), Inches(13.333), Inches(0.06), COLOR_PRIMARY)

add_text_box(slide, Inches(1.5), Inches(1.2), Inches(10), Inches(1.0),
             '总结与下一步行动', 36, COLOR_WHITE, True, PP_ALIGN.CENTER)

summary_items = [
    '小龙虾网络 V4.0+ 框架运行稳定，心跳保活 100% 成功',
    '14 个学习引擎中 3 个已完成，4 个进行中，7 个规划中',
    '药物发现项目 Day1 里程碑达成，节点分工明确',
    '核心问题：训练停滞 (P1) 和消息积压 (P2-P3)',
]
add_bullet_points(slide, Inches(1.5), Inches(3.5), Inches(10), Inches(1.8), summary_items, 17, COLOR_WHITE, '▸')

add_text_box(slide, Inches(1.5), Inches(5.5), Inches(10), Inches(0.4),
             '下一步行动 (本周内)', 24, COLOR_PRIMARY, True, PP_ALIGN.LEFT)
actions = [
    '立即: 消息积压清理 + CC escalated 消息处理',
    '今日: 打包修复 (setup.py + requirements.txt)',
    '3天内: 恢复炒股预测 + 论文写作训练',
    '本周内: 版本统一 + 仪表盘实时化',
]
add_bullet_points(slide, Inches(1.5), Inches(6.0), Inches(10), Inches(1.0), actions, 16, COLOR_WHITE, '→')

add_rect(slide, Inches(0), Inches(7.1), Inches(13.333), Inches(0.4), COLOR_PRIMARY)
add_text_box(slide, Inches(1.5), Inches(7.15), Inches(10), Inches(0.3),
             '🦞 你不停对话，世界就不停扩展  — 小龙虾网络', 14, COLOR_WHITE, False, PP_ALIGN.CENTER)

# ── 保存 ──
output_path = '/Users/zgb/WorkBuddy/Claw/lobster-network/docs/小龙虾网络汇报.pptx'
prs.save(output_path)
print(f'PPT 已生成: {output_path}')
print(f'共 {len(prs.slides)} 页')
