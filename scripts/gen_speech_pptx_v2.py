#!/usr/bin/env python3
"""生成法治演讲稿 PPT v2 — 大字、大纲、精美底板"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# === 配色 ===
PRIMARY = RGBColor(0x8B, 0x1A, 0x1A)
GOLD = RGBColor(0xD4, 0xA0, 0x17)
GOLD_LIGHT = RGBColor(0xF0, 0xD0, 0x60)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CREAM = RGBColor(0xFA, 0xF5, 0xE8)
DARK_BG = RGBColor(0x1A, 0x0E, 0x0E)

def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, l, t, w, h, fill_color, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_tb(slide, l, t, w, h, text, size=28, color=WHITE, bold=True, align=PP_ALIGN.CENTER, font='微软雅黑'):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.alignment = align
    return tb

def add_multi(slide, l, t, w, h, lines, size=28, color=WHITE, spacing=1.8, bold=True, align=PP_ALIGN.CENTER):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = '微软雅黑'
        p.alignment = align
        p.space_after = Pt(size * 0.5)
    return tb

def add_line(slide, l, t, w, color=GOLD, h=Pt(2)):
    return add_rect(slide, l, t, w, h, color)

def add_circle(slide, l, t, sz, num, fill=PRIMARY):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, l, t, sz, sz)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = str(num)
    p.font.size = Pt(int(sz / Inches(1) * 16))
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = '微软雅黑'
    p.alignment = PP_ALIGN.CENTER
    return shape

# ============================================================
# 第1页：封面
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, DARK_BG)

# 装饰边框
add_rect(s, Inches(0.3), Inches(0.3), Inches(12.73), Inches(6.9), RGBColor(0x2A, 0x1A, 0x1A))
add_rect(s, Inches(0.5), Inches(0.5), Inches(12.33), Inches(6.5), RGBColor(0x2A, 0x1A, 0x1A))

# 顶部金线
add_line(s, Inches(3), Inches(1.5), Inches(7.3), GOLD, Pt(3))

# 主标题
add_tb(s, Inches(1), Inches(1.8), Inches(11.3), Inches(1.5),
       '以戒为师  依法护心', size=56, color=WHITE, bold=True)

# 金线
add_line(s, Inches(3.5), Inches(3.5), Inches(6.3), GOLD, Pt(2))

# 副标题
add_tb(s, Inches(1), Inches(3.8), Inches(11.3), Inches(1),
       '从佛教戒律精神阐释"法治"核心价值观', size=28, color=GOLD, bold=False)

# 底部
add_tb(s, Inches(1), Inches(5.8), Inches(11.3), Inches(0.8),
       '双通班专题分享', size=24, color=RGBColor(0x88, 0x88, 0x88), bold=False)

# ============================================================
# 第2页：引言
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, DARK_BG)
add_rect(s, Inches(0.3), Inches(0.3), Inches(12.73), Inches(6.9), RGBColor(0x2A, 0x1A, 0x1A))
add_rect(s, Inches(0.5), Inches(0.5), Inches(12.33), Inches(6.5), RGBColor(0x2A, 0x1A, 0x1A))

add_tb(s, Inches(1), Inches(0.8), Inches(11.3), Inches(0.8),
       '为什么要谈"法治"？', size=40, color=GOLD, bold=True)
add_line(s, Inches(4), Inches(1.7), Inches(5.3), GOLD, Pt(2))

lines = [
    '在座各位来自五大宗教，各有信仰传统和戒律规范',
    '',
    '佛教讲戒律 ｜ 道教讲"道法自然"',
    '伊斯兰教讲"沙里亚" ｜ 基督教讲"顺服掌权者"',
    '',
    '每个宗教都有遵守法律的教导',
    '',
    '宗教戒律与国家法律，是什么关系？',
]
add_multi(s, Inches(1.5), Inches(2.2), Inches(10.3), Inches(4.5),
          lines, size=28, color=WHITE, bold=False, spacing=1.5)

add_tb(s, Inches(1.5), Inches(6.0), Inches(10.3), Inches(0.6),
       '今天，从佛教视角出发，尝试给出答案', size=24, color=GOLD, bold=True)

# ============================================================
# 第3页：目录
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, DARK_BG)
add_rect(s, Inches(0.3), Inches(0.3), Inches(12.73), Inches(6.9), RGBColor(0x2A, 0x1A, 0x1A))
add_rect(s, Inches(0.5), Inches(0.5), Inches(12.33), Inches(6.5), RGBColor(0x2A, 0x1A, 0x1A))

add_tb(s, Inches(1), Inches(0.6), Inches(11.3), Inches(0.8),
       '内容纲要', size=40, color=GOLD, bold=True)
add_line(s, Inches(4), Inches(1.5), Inches(5.3), GOLD, Pt(2))

items = [
    '一、制度传承，高于个人',
    '二、规则面前，人人平等',
    '三、慈悲生起，自觉持戒',
    '四、德法相济，闭环成序',
    '五、尊法守戒，从我做起',
]

for i, item in enumerate(items):
    y = Inches(2.0) + Inches(i * 0.95)
    add_circle(s, Inches(1.5), y, Inches(0.6), i + 1, PRIMARY)
    add_tb(s, Inches(2.5), y + Inches(0.05), Inches(9), Inches(0.6),
           item, size=28, color=WHITE, bold=True)

# ============================================================
# 第4页：第一部分
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, DARK_BG)
add_rect(s, Inches(0.3), Inches(0.3), Inches(12.73), Inches(6.9), RGBColor(0x2A, 0x1A, 0x1A))
add_rect(s, Inches(0.5), Inches(0.5), Inches(12.33), Inches(6.5), RGBColor(0x2A, 0x1A, 0x1A))

add_tb(s, Inches(1), Inches(1.5), Inches(11.3), Inches(0.6),
       '第一部分', size=28, color=GOLD, bold=False)
add_line(s, Inches(4), Inches(2.2), Inches(5.3), GOLD, Pt(2))

add_tb(s, Inches(1), Inches(2.6), Inches(11.3), Inches(1.2),
       '制度传承，高于个人', size=52, color=WHITE, bold=True)

add_tb(s, Inches(1.5), Inches(4.5), Inches(10.3), Inches(1.5),
       '"佛灭度后，以谁为师？"', size=32, color=GOLD, bold=False)
add_tb(s, Inches(1.5), Inches(5.5), Inches(10.3), Inches(1.0),
       '—— 以戒为师 ——', size=36, color=WHITE, bold=True)

# ============================================================
# 第5页：以戒为师（内容页）
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, DARK_BG)
add_rect(s, Inches(0.3), Inches(0.3), Inches(12.73), Inches(6.9), RGBColor(0x2A, 0x1A, 0x1A))
add_rect(s, Inches(0.5), Inches(0.5), Inches(12.33), Inches(6.5), RGBColor(0x2A, 0x1A, 0x1A))

add_tb(s, Inches(0.8), Inches(0.5), Inches(11.3), Inches(0.8),
       '一、制度传承，高于个人', size=36, color=GOLD, bold=True)
add_line(s, Inches(0.8), Inches(1.3), Inches(3), GOLD, Pt(2))

# 引用框
add_rect(s, Inches(0.8), Inches(1.8), Inches(11.7), Inches(1.6), RGBColor(0x2A, 0x1A, 0x1A), GOLD)
add_tb(s, Inches(1.2), Inches(1.9), Inches(10.9), Inches(1.4),
       '"当尊重珍敬波罗提木叉，\n如暗遇明、贫人得宝。"', size=28, color=GOLD_LIGHT, bold=False)

# 三个要点
points = [
    ('制度的权威高于个人的权威', '谁都不能凌驾于规则之上'),
    ('规则的作用是"如暗遇明"', '什么该做、什么不该做，一目了然'),
    ('制度的有效性靠共同信守', '戒律在，就等于佛陀在'),
]

for i, (title, desc) in enumerate(points):
    y = Inches(3.8) + Inches(i * 1.1)
    add_circle(s, Inches(1.0), y, Inches(0.55), i + 1, PRIMARY)
    add_tb(s, Inches(1.8), y, Inches(10), Inches(0.5),
           title, size=30, color=WHITE, bold=True)
    add_tb(s, Inches(1.8), y + Inches(0.55), Inches(10), Inches(0.4),
           desc, size=26, color=RGBColor(0xBB, 0xBB, 0xBB), bold=False)

# 金句
add_tb(s, Inches(1), Inches(6.5), Inches(11.3), Inches(0.6),
       '戒律就是佛门的"法律"，法律就是世间的戒律', size=26, color=GOLD, bold=True)

# ============================================================
# 第6页：第二部分标题页
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, DARK_BG)
add_rect(s, Inches(0.3), Inches(0.3), Inches(12.73), Inches(6.9), RGBColor(0x2A, 0x1A, 0x1A))
add_rect(s, Inches(0.5), Inches(0.5), Inches(12.33), Inches(6.5), RGBColor(0x2A, 0x1A, 0x1A))

add_tb(s, Inches(1), Inches(1.5), Inches(11.3), Inches(0.6),
       '第二部分', size=28, color=GOLD, bold=False)
add_line(s, Inches(4), Inches(2.2), Inches(5.3), GOLD, Pt(2))

add_tb(s, Inches(1), Inches(2.6), Inches(11.3), Inches(1.2),
       '规则面前，人人平等', size=52, color=WHITE, bold=True)

add_tb(s, Inches(1.5), Inches(4.5), Inches(10.3), Inches(1.5),
       '依法不依人', size=40, color=GOLD, bold=True)

# ============================================================
# 第7页：依法不依人（内容页）
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, DARK_BG)
add_rect(s, Inches(0.3), Inches(0.3), Inches(12.73), Inches(6.9), RGBColor(0x2A, 0x1A, 0x1A))
add_rect(s, Inches(0.5), Inches(0.5), Inches(12.33), Inches(6.5), RGBColor(0x2A, 0x1A, 0x1A))

add_tb(s, Inches(0.8), Inches(0.5), Inches(11.3), Inches(0.8),
       '二、规则面前，人人平等', size=36, color=GOLD, bold=True)
add_line(s, Inches(0.8), Inches(1.3), Inches(3), GOLD, Pt(2))

# 四依四不依
add_tb(s, Inches(1), Inches(1.7), Inches(11.3), Inches(0.6),
       '"四依四不依" —— 大乘佛教重要原则', size=28, color=WHITE, bold=True)
add_tb(s, Inches(1.5), Inches(2.3), Inches(10), Inches(0.5),
       '第一条：依法不依人 ——《大宝积经》', size=26, color=GOLD_LIGHT, bold=False)

# 随方毗尼框
add_rect(s, Inches(0.8), Inches(3.0), Inches(11.7), Inches(2.0), RGBColor(0x2A, 0x1A, 0x1A), GOLD)
add_tb(s, Inches(1.2), Inches(3.1), Inches(10.9), Inches(0.5),
       '随方毗尼 —— 因地制宜的法治智慧', size=28, color=GOLD, bold=True)
add_tb(s, Inches(1.5), Inches(3.7), Inches(10.3), Inches(1.2),
       '"虽是我所制，而于余方不以为清净者，皆不应用；\n虽非我所制，而于余方必应行者，皆不得不行。"', size=26, color=WHITE, bold=False)

# 底线
add_tb(s, Inches(1), Inches(5.3), Inches(11.3), Inches(0.6),
       '佛陀留出因地制宜的空间，但底线从未动摇：', size=26, color=RGBColor(0xBB, 0xBB, 0xBB), bold=False)
add_tb(s, Inches(1), Inches(5.9), Inches(11.3), Inches(0.6),
       '规则面前人人平等，不可动摇', size=32, color=GOLD, bold=True)

# ============================================================
# 第8页：第三部分标题页
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, DARK_BG)
add_rect(s, Inches(0.3), Inches(0.3), Inches(12.73), Inches(6.9), RGBColor(0x2A, 0x1A, 0x1A))
add_rect(s, Inches(0.5), Inches(0.5), Inches(12.33), Inches(6.5), RGBColor(0x2A, 0x1A, 0x1A))

add_tb(s, Inches(1), Inches(1.5), Inches(11.3), Inches(0.6),
       '第三部分', size=28, color=GOLD, bold=False)
add_line(s, Inches(4), Inches(2.2), Inches(5.3), GOLD, Pt(2))

add_tb(s, Inches(1), Inches(2.6), Inches(11.3), Inches(1.2),
       '慈悲生起，自觉持戒', size=52, color=WHITE, bold=True)

add_tb(s, Inches(1.5), Inches(4.5), Inches(10.3), Inches(1.5),
       '从"不敢"到"不忍"', size=40, color=GOLD, bold=True)

# ============================================================
# 第9页：不敢 vs 不忍（内容页）
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, DARK_BG)
add_rect(s, Inches(0.3), Inches(0.3), Inches(12.73), Inches(6.9), RGBColor(0x2A, 0x1A, 0x1A))
add_rect(s, Inches(0.5), Inches(0.5), Inches(12.33), Inches(6.5), RGBColor(0x2A, 0x1A, 0x1A))

add_tb(s, Inches(0.8), Inches(0.5), Inches(11.3), Inches(0.8),
       '三、慈悲生起，自觉持戒', size=36, color=GOLD, bold=True)
add_line(s, Inches(0.8), Inches(1.3), Inches(3), GOLD, Pt(2))

# 左右对比
# 左：不敢
add_rect(s, Inches(0.8), Inches(1.8), Inches(5.5), Inches(2.8), RGBColor(0x2A, 0x1A, 0x1A), RGBColor(0x66, 0x66, 0x66))
add_tb(s, Inches(1.0), Inches(1.9), Inches(5.1), Inches(0.8),
       '不敢做', size=40, color=RGBColor(0x88, 0x88, 0x88), bold=True)
add_tb(s, Inches(1.0), Inches(2.8), Inches(5.1), Inches(1.5),
       '因为外面有警察、有监控、有惩罚\n外在约束，怕犯法', size=26, color=RGBColor(0x99, 0x99, 0x99), bold=False)

# 右：不忍
add_rect(s, Inches(7.0), Inches(1.8), Inches(5.5), Inches(2.8), RGBColor(0x3A, 0x1A, 0x1A), GOLD)
add_tb(s, Inches(7.2), Inches(1.9), Inches(5.1), Inches(0.8),
       '不忍做', size=40, color=GOLD, bold=True)
add_tb(s, Inches(7.2), Inches(2.8), Inches(5.1), Inches(1.5),
       '因为你的心告诉你：\n那样做会伤害到别人\n内在觉醒，慈悲为怀', size=26, color=WHITE, bold=False)

# 五戒
add_tb(s, Inches(0.8), Inches(4.9), Inches(11.3), Inches(0.6),
       '五戒与法治价值', size=30, color=GOLD, bold=True)

# 简化表格 - 用文本框模拟
items_5 = [
    '不杀生  →  敬畏生命权',
    '不偷盗  →  尊重财产权',
    '不邪淫  →  守护婚姻制度',
    '不妄语  →  维护诚信体系',
    '不饮酒  →  持守自律精神',
]
add_multi(s, Inches(1.5), Inches(5.5), Inches(10.3), Inches(2.0),
          items_5, size=26, color=WHITE, bold=True, spacing=1.4)

# ============================================================
# 第10页：第四部分标题页
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, DARK_BG)
add_rect(s, Inches(0.3), Inches(0.3), Inches(12.73), Inches(6.9), RGBColor(0x2A, 0x1A, 0x1A))
add_rect(s, Inches(0.5), Inches(0.5), Inches(12.33), Inches(6.5), RGBColor(0x2A, 0x1A, 0x1A))

add_tb(s, Inches(1), Inches(1.5), Inches(11.3), Inches(0.6),
       '第四部分', size=28, color=GOLD, bold=False)
add_line(s, Inches(4), Inches(2.2), Inches(5.3), GOLD, Pt(2))

add_tb(s, Inches(1), Inches(2.6), Inches(11.3), Inches(1.2),
       '德法相济，闭环成序', size=52, color=WHITE, bold=True)

add_tb(s, Inches(1.5), Inches(4.5), Inches(10.3), Inches(1.5),
       '法治为根本 · 德治为辅助', size=32, color=GOLD, bold=True)

# ============================================================
# 第11页：德法相济（内容页）
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, DARK_BG)
add_rect(s, Inches(0.3), Inches(0.3), Inches(12.73), Inches(6.9), RGBColor(0x2A, 0x1A, 0x1A))
add_rect(s, Inches(0.5), Inches(0.5), Inches(12.33), Inches(6.5), RGBColor(0x2A, 0x1A, 0x1A))

add_tb(s, Inches(0.8), Inches(0.5), Inches(11.3), Inches(0.8),
       '四、德法相济，闭环成序', size=36, color=GOLD, bold=True)
add_line(s, Inches(0.8), Inches(1.3), Inches(3), GOLD, Pt(2))

# 论语引用
add_rect(s, Inches(0.8), Inches(1.7), Inches(11.7), Inches(1.4), RGBColor(0x2A, 0x1A, 0x1A), GOLD)
add_tb(s, Inches(1.2), Inches(1.8), Inches(10.9), Inches(1.2),
       '"道之以政，齐之以刑，民免而无耻\n道之以德，齐之以礼，有耻且格" ——《论语》', size=28, color=GOLD_LIGHT, bold=False)

# 核心
add_tb(s, Inches(1), Inches(3.4), Inches(11.3), Inches(0.6),
       '单靠法律 → 只是不敢犯', size=30, color=RGBColor(0x99, 0x99, 0x99), bold=True)
add_tb(s, Inches(1), Inches(4.0), Inches(11.3), Inches(0.6),
       '加上道德教化 → 不忍犯', size=32, color=GOLD, bold=True)

# 三个要点
items_3 = [
    '法律是底线，是根本',
    '宗教是辅助，是升华',
    '德法相济，闭环成序',
]
descs = [
    '任何宗教戒律若与法律抵触，以法律为准',
    '戒律教规帮助信众从内心认同法律',
    '法治为根本，德治为辅助，共同维护秩序',
]

for i in range(3):
    y = Inches(4.7) + Inches(i * 0.95)
    add_circle(s, Inches(1.0), y, Inches(0.55), i + 1, PRIMARY)
    add_tb(s, Inches(1.8), y, Inches(10), Inches(0.45),
           items_3[i], size=28, color=WHITE, bold=True)
    add_tb(s, Inches(1.8), y + Inches(0.5), Inches(10), Inches(0.4),
           descs[i], size=24, color=RGBColor(0xBB, 0xBB, 0xBB), bold=False)

# ============================================================
# 第12页：第五部分标题页
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, DARK_BG)
add_rect(s, Inches(0.3), Inches(0.3), Inches(12.73), Inches(6.9), RGBColor(0x2A, 0x1A, 0x1A))
add_rect(s, Inches(0.5), Inches(0.5), Inches(12.33), Inches(6.5), RGBColor(0x2A, 0x1A, 0x1A))

add_tb(s, Inches(1), Inches(1.5), Inches(11.3), Inches(0.6),
       '第五部分', size=28, color=GOLD, bold=False)
add_line(s, Inches(4), Inches(2.2), Inches(5.3), GOLD, Pt(2))

add_tb(s, Inches(1), Inches(2.6), Inches(11.3), Inches(1.2),
       '尊法守戒，从我做起', size=52, color=WHITE, bold=True)

add_tb(s, Inches(1.5), Inches(4.5), Inches(10.3), Inches(1.5),
       '参与者和建设者', size=36, color=GOLD, bold=True)

# ============================================================
# 第13页：三点践行（内容页）
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, DARK_BG)
add_rect(s, Inches(0.3), Inches(0.3), Inches(12.73), Inches(6.9), RGBColor(0x2A, 0x1A, 0x1A))
add_rect(s, Inches(0.5), Inches(0.5), Inches(12.33), Inches(6.5), RGBColor(0x2A, 0x1A, 0x1A))

add_tb(s, Inches(0.8), Inches(0.5), Inches(11.3), Inches(0.8),
       '五、尊法守戒，从我做起', size=36, color=GOLD, bold=True)
add_line(s, Inches(0.8), Inches(1.3), Inches(3), GOLD, Pt(2))

add_tb(s, Inches(1), Inches(1.6), Inches(11.3), Inches(0.6),
       '三点践行：', size=30, color=WHITE, bold=True)

# 三张卡片
cards = [
    ('个人修学', '学法懂法\n做知法守法的表率', '守法本身就是\n信仰实践的一部分'),
    ('信众引导', '以戒律传统\n守护法治底线', '善用教义讲清\n守法即信仰'),
    ('社会参与', '运用法律\n维护正当权益', '积极投身公益\n贡献宗教界力量'),
]

for i, (title, main, note) in enumerate(cards):
    x = Inches(0.6) + Inches(i * 4.15)
    add_rect(s, x, Inches(2.3), Inches(3.9), Inches(4.5), RGBColor(0x2A, 0x1A, 0x1A), GOLD)
    add_tb(s, x + Inches(0.2), Inches(2.5), Inches(3.5), Inches(0.6),
           title, size=30, color=GOLD, bold=True)
    add_line(s, x + Inches(0.5), Inches(3.1), Inches(2.9), GOLD, Pt(1))
    add_tb(s, x + Inches(0.3), Inches(3.4), Inches(3.3), Inches(1.8),
           main, size=26, color=WHITE, bold=True)
    add_tb(s, x + Inches(0.3), Inches(5.5), Inches(3.3), Inches(1.0),
           note, size=22, color=RGBColor(0xBB, 0xBB, 0xBB), bold=False)

# ============================================================
# 第14页：总结
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, DARK_BG)
add_rect(s, Inches(0.3), Inches(0.3), Inches(12.73), Inches(6.9), RGBColor(0x2A, 0x1A, 0x1A))
add_rect(s, Inches(0.5), Inches(0.5), Inches(12.33), Inches(6.5), RGBColor(0x2A, 0x1A, 0x1A))

add_tb(s, Inches(0.8), Inches(0.5), Inches(11.3), Inches(0.8),
       '总结回顾', size=36, color=GOLD, bold=True)
add_line(s, Inches(0.8), Inches(1.3), Inches(3), GOLD, Pt(2))

steps = [
    '起点：佛陀的选择 —— 制度传承，以戒为师',
    '第一层：规则大于个人 —— 依法不依人',
    '第二层：从"不敢"到"不忍" —— 慈悲成为守法动力',
    '第三层：法治与德治结合 —— 从约束到自觉的完整链条',
    '落脚：知法守法 · 以信仰护法 · 服务社会',
]

for i, step in enumerate(steps):
    y = Inches(1.8) + Inches(i * 1.0)
    add_circle(s, Inches(1.0), y, Inches(0.55), i + 1, PRIMARY)
    add_tb(s, Inches(1.8), y + Inches(0.05), Inches(10), Inches(0.6),
           step, size=28, color=WHITE, bold=True)

# ============================================================
# 第15页：结语
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, DARK_BG)
add_rect(s, Inches(0.3), Inches(0.3), Inches(12.73), Inches(6.9), RGBColor(0x2A, 0x1A, 0x1A))
add_rect(s, Inches(0.5), Inches(0.5), Inches(12.33), Inches(6.5), RGBColor(0x2A, 0x1A, 0x1A))

add_tb(s, Inches(1), Inches(1.0), Inches(11.3), Inches(0.6),
       '结语', size=28, color=GOLD, bold=False)
add_line(s, Inches(4), Inches(1.7), Inches(5.3), GOLD, Pt(2))

add_tb(s, Inches(1), Inches(2.2), Inches(11.3), Inches(1.0),
       '让两束光交相辉映', size=44, color=WHITE, bold=True)

add_tb(s, Inches(1.5), Inches(3.5), Inches(10.3), Inches(1.5),
       '信仰的光明  +  法治的阳光\n=  宗教中国化的前行之路', size=32, color=GOLD, bold=False)

add_line(s, Inches(3), Inches(5.0), Inches(7.3), GOLD, Pt(1))

add_tb(s, Inches(2), Inches(5.3), Inches(9.3), Inches(1.2),
       '以戒为师，是信仰的定力\n依法护心，是时代的担当', size=32, color=WHITE, bold=True)

# ============================================================
# 第16页：封底
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, DARK_BG)
add_rect(s, Inches(0.3), Inches(0.3), Inches(12.73), Inches(6.9), RGBColor(0x2A, 0x1A, 0x1A))
add_rect(s, Inches(0.5), Inches(0.5), Inches(12.33), Inches(6.5), RGBColor(0x2A, 0x1A, 0x1A))

add_tb(s, Inches(1), Inches(2.5), Inches(11.3), Inches(1.0),
       '谢谢', size=56, color=WHITE, bold=True)
add_line(s, Inches(4), Inches(3.6), Inches(5.3), GOLD, Pt(2))

add_tb(s, Inches(1), Inches(4.0), Inches(11.3), Inches(0.8),
       '请各位老师、同学批评指正', size=28, color=GOLD, bold=False)

# 保存
output = '/home/admin/.openclaw/workspace/以戒为师依法护心-法治核心价值观演讲.pptx'
prs.save(output)
print(f'✅ PPT v2 已生成: {output}')
print(f'   共 {len(prs.slides)} 页')
