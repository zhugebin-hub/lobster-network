#!/usr/bin/env python3
"""
更新小龙虾网络项目汇报 PPT 到 v0.4.1
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import sys

def set_cell_text(cell, text, font_size=12, bold=False, color=None):
    """设置单元格文本"""
    cell.text = ""
    p = cell.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    if color:
        p.font.color.rgb = color
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE

def add_text_box(slide, left, top, width, height, text, font_size=18, bold=False, color=None, alignment=PP_ALIGN.LEFT):
    """添加文本框"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    if color:
        p.font.color.rgb = color
    p.alignment = alignment
    return txBox

def add_bullet_list(slide, left, top, width, height, items, font_size=14, color=None):
    """添加项目符号列表"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.space_after = Pt(6)
        if color:
            p.font.color.rgb = color
        p.level = 0
    
    return txBox

def create_updated_ppt():
    """创建更新后的 PPT"""
    prs = Presentation('/home/admin/.openclaw/media/inbound/c242a8a4-38d5-4ad0-bac6-39cc3691ff15.pptx')
    
    # 更新第1页 - 封面
    slide1 = prs.slides[0]
    for shape in slide1.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                if "v0.3.0" in para.text:
                    para.text = para.text.replace("v0.3.0", "v0.4.1 整合版")
                if "2026年6月" in para.text:
                    para.text = para.text.replace("2026年6月", "2026年6月24日")
    
    # 更新第5页 - 四层架构 → 五层架构
    slide5 = prs.slides[4]
    # 在框架层和基础设施层之间插入可靠通信层说明
    # 由于PPT结构复杂，我们添加一个备注框
    add_text_box(
        slide5, Inches(0.5), Inches(5.5), Inches(9), Inches(0.5),
        "【v0.4.1 新增】可靠通信层：节点注册中心 + 可靠消息传递 + 多通道故障切换",
        font_size=12, bold=True, color=RGBColor(0x00, 0x70, 0xC0)
    )
    
    # 更新第6页 - 7个核心模块 → 10个核心模块
    slide6 = prs.slides[5]
    for shape in slide6.shapes:
        if shape.has_table:
            table = shape.table
            # 添加3行新模块
            for module in [
                ("节点注册中心", "NodeRegistry", "注册/发现/心跳/健康检查/持久化"),
                ("可靠消息传递", "Messenger", "ACK/NACK/重试/故障切换/持久化"),
                ("集成层", "LobsterNetworkWithRegistry", "注册中心+消息+网络整合"),
            ]:
                row = table.rows.add()
                for i, text in enumerate(module):
                    set_cell_text(row.cells[i], text, font_size=11)
    
    # 更新第11页 - 版本历程
    slide11 = prs.slides[10]
    for shape in slide11.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                if "v0.4.0" in para.text:
                    para.text = para.text.replace("v0.4.0", "v0.4.0/v0.4.1")
                if "计划中" in para.text:
                    para.text = para.text.replace("计划中", "✅ 已完成")
    
    # 更新第12页 - 关键数据
    slide12 = prs.slides[11]
    # 添加新数据
    add_text_box(
        slide12, Inches(0.5), Inches(5.0), Inches(9), Inches(0.8),
        "📊 v0.4.1 新增数据：\n• 测试用例：62个全部通过（虾尔37 + 诸葛马25）\n• 传输通道：4种（NFS/SSH/HTTP/File）\n• 部署方式：一键部署 + 一键回滚",
        font_size=14, bold=True, color=RGBColor(0x00, 0x70, 0xC0)
    )
    
    # 更新第13页 - 路线图
    slide13 = prs.slides[12]
    for shape in slide13.shapes:
        if shape.has_table:
            table = shape.table
            # 更新现有行
            for row in table.rows:
                cells = [cell.text for cell in row.cells]
                if "v0.4.0" in cells[0]:
                    set_cell_text(row.cells[0], "v0.4.0/v0.4.1", font_size=11)
                    set_cell_text(row.cells[2], "✅ 已完成", font_size=11, color=RGBColor(0x00, 0xB0, 0x50))
                elif "v0.5.0" in cells[0]:
                    set_cell_text(row.cells[0], "v0.4.2", font_size=11)
                    set_cell_text(row.cells[1], "安全增强 + 监控告警 + 性能优化", font_size=11)
                    set_cell_text(row.cells[2], "🔄 规划中", font_size=11, color=RGBColor(0xFF, 0xC0, 0x00))
            
            # 添加新版本
            for version in [
                ("v0.5.0", "分布式架构 + 跨域协作", "📋 计划中"),
                ("v1.0.0", "正式发布", "🎯 目标"),
            ]:
                row = table.rows.add()
                for i, text in enumerate(version):
                    set_cell_text(row.cells[i], text, font_size=11)
    
    # 保存
    output_path = '/tmp/lobster-network-项目汇报-v0.4.1.pptx'
    prs.save(output_path)
    print(f"✅ 更新完成: {output_path}")
    return output_path

if __name__ == "__main__":
    create_updated_ppt()
