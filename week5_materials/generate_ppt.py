#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
第 5 周 PPT 生成脚本
数字人文与小龙虾课程
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def create_week5_ppt():
    # 创建演示文稿
    prs = Presentation()
    
    # 设置幻灯片大小为 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # 定义颜色
    BLUE = RGBColor(0, 82, 155)  # 浙商大蓝
    ORANGE = RGBColor(255, 140, 0)  # 小龙虾橙
    DARK_GRAY = RGBColor(64, 64, 64)
    LIGHT_GRAY = RGBColor(240, 240, 240)
    
    def add_title_slide(title, subtitle):
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        title_shape = slide.shapes.title
        title_shape.text = title
        title_shape.text_frame.paragraphs[0].font.size = Pt(44)
        title_shape.text_frame.paragraphs[0].font.bold = True
        title_shape.text_frame.paragraphs[0].font.color.rgb = BLUE
        
        subtitle_shape = slide.placeholders[1]
        subtitle_shape.text = subtitle
        subtitle_shape.text_frame.paragraphs[0].font.size = Pt(24)
        subtitle_shape.text_frame.paragraphs[0].font.color.rgb = DARK_GRAY
        return slide
    
    def add_content_slide(title, content_items):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title_shape = slide.shapes.title
        title_shape.text = title
        title_shape.text_frame.paragraphs[0].font.size = Pt(36)
        title_shape.text_frame.paragraphs[0].font.bold = True
        title_shape.text_frame.paragraphs[0].font.color.rgb = BLUE
        
        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        tf.clear()
        
        for i, item in enumerate(content_items):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = item
            p.font.size = Pt(22)
            p.space_after = Pt(14)
            if item.startswith("  "):
                p.level = 1
                p.font.size = Pt(20)
        
        return slide
    
    def add_bullet_slide(title, bullets):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title_shape = slide.shapes.title
        title_shape.text = title
        title_shape.text_frame.paragraphs[0].font.size = Pt(36)
        title_shape.text_frame.paragraphs[0].font.bold = True
        title_shape.text_frame.paragraphs[0].font.color.rgb = BLUE
        
        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        tf.clear()
        
        for i, bullet in enumerate(bullets):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            
            if isinstance(bullet, tuple):
                text, level = bullet
                p.text = text
                p.level = level
            else:
                p.text = bullet
                p.level = 0
            
            p.font.size = Pt(20)
            p.space_after = Pt(10)
        
        return slide
    
    def add_table_slide(title, headers, rows):
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        
        # 添加标题
        title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(1))
        tf = title_shape.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = BLUE
        
        # 添加表格
        rows_count = len(rows) + 1
        cols_count = len(headers)
        table = slide.shapes.add_table(
            rows_count, cols_count,
            Inches(0.5), Inches(1.5),
            Inches(12.333), Inches(5.5)
        ).table
        
        # 设置列宽
        for i in range(cols_count):
            table.columns[i].width = Inches(12.333 / cols_count)
        
        # 填充表头
        for i, header in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = header
            cell.fill.solid()
            cell.fill.fore_color.rgb = BLUE
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(18)
                paragraph.font.bold = True
                paragraph.font.color.rgb = RGBColor(255, 255, 255)
                paragraph.alignment = PP_ALIGN.CENTER
        
        # 填充数据行
        for row_idx, row_data in enumerate(rows, 1):
            for col_idx, cell_data in enumerate(row_data):
                cell = table.cell(row_idx, col_idx)
                cell.text = cell_data
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.size = Pt(16)
                    paragraph.alignment = PP_ALIGN.LEFT
        
        return slide
    
    # ========== 第一部分：课程导入 ==========
    
    # Slide 1: 封面页
    slide = add_title_slide(
        "第 5 周 互动与生态养成",
        "数字人文与小龙虾\n浙江工商大学通识课程\n2026 年 3 月 30 日"
    )
    
    # Slide 2: 本周学习目标
    slide = add_bullet_slide("本周学习目标", [
        "📚 知识目标：理解人机协作的深层逻辑与边界",
        "🛠️ 技能目标：掌握高阶提问、追问、验证技巧",
        "💡 素养目标：建立批判性使用 AI 的思维习惯",
        "🤝 情感目标：与 AI 建立健康的'共生关系'认知"
    ])
    
    # Slide 3: 前 4 周回顾
    slide = add_bullet_slide("前 4 周回顾", [
        "第 1 周：数字人文是什么（学科认知）",
        "第 2 周：认识你的小龙虾（AI 基础）",
        "第 3 周：精准提问的艺术（结构化表达）",
        "第 4 周：验证与边界（批判性思维）",
        "",
        "本周：从'怎么用'到'怎么养'"
    ])
    
    # Slide 4: 小龙虾三部曲回顾
    slide = add_bullet_slide("小龙虾三部曲核心要点", [
        "🦞 认识小龙虾",
        ("能力与边界：能做什么、不能做什么", 1),
        ("栖息地：阿里云服务器", 1),
        ("数量：每人 1 只", 1),
        "",
        "🎯 精准提问",
        ("结构化表达：背景 + 任务 + 要求 + 格式", 1),
        ("好问题 = 好答案", 1),
        "",
        "✅ 验证答案",
        ("AI 可能出错，需要验证", 1),
        ("学术诚信红线", 1)
    ])
    
    # Slide 5: 快问快答
    slide = add_bullet_slide("快问快答（互动环节）", [
        "❓ 问题 1：小龙虾的'栖息地'在哪里？",
        "❓ 问题 2：每人建议养几只小龙虾？",
        "❓ 问题 3：AI 能帮你考试作弊吗？",
        "❓ 问题 4：AI 出错了怎么办？",
        "",
        "💡 答案稍后揭晓..."
    ])
    
    # ========== 第二部分：商科 AI 应用场景 ==========
    
    # Slide 6: 场景导入
    slide = add_title_slide(
        "AI 在商科中的 5 大应用场景",
        "你的专业能用 AI 做什么？"
    )
    
    # Slide 7: 场景 1
    slide = add_bullet_slide("场景 1：年报文本分析", [
        "📊 适用专业：会计学、财务管理",
        "",
        "💼 应用内容：",
        ("用 AI 分析企业年报情感倾向", 1),
        ("提取管理层讨论的关键信息", 1),
        "",
        "❓ 示例问题：",
        ("'分析某公司年报中管理层讨论的情感变化'", 1),
        ("'提取年报中的风险因素关键词'", 1)
    ])
    
    # Slide 8: 场景 2
    slide = add_bullet_slide("场景 2：品牌舆情监测", [
        "📱 适用专业：市场营销、电子商务",
        "",
        "💼 应用内容：",
        ("社交媒体评论情感分析", 1),
        ("品牌口碑追踪", 1),
        "",
        "❓ 示例问题：",
        ("'分析某品牌微博评论的情感倾向'", 1),
        ("'提取消费者最关注的产品特征'", 1)
    ])
    
    # Slide 9: 场景 3
    slide = add_bullet_slide("场景 3：商业史梳理", [
        "🏢 适用专业：工商管理、国际贸易",
        "",
        "💼 应用内容：",
        ("企业发展历程知识图谱", 1),
        ("商业事件脉络整理", 1),
        "",
        "❓ 示例问题：",
        ("'梳理某企业 10 年发展关键事件'", 1),
        ("'构建浙商代表人物关系网络'", 1)
    ])
    
    # Slide 10: 场景 4
    slide = add_bullet_slide("场景 4：消费者评论挖掘", [
        "🛒 适用专业：电子商务、市场营销",
        "",
        "💼 应用内容：",
        ("电商评论关键词提取", 1),
        ("用户反馈分析", 1),
        "",
        "❓ 示例问题：",
        ("'分析某产品差评的主要原因'", 1),
        ("'提取影响购买决策的关键因素'", 1)
    ])
    
    # Slide 11: 场景 5
    slide = add_bullet_slide("场景 5：财经新闻追踪", [
        "📰 适用专业：金融学、经济学",
        "",
        "💼 应用内容：",
        ("行业新闻趋势可视化", 1),
        ("政策影响分析", 1),
        "",
        "❓ 示例问题：",
        ("'分析某行业近半年新闻热点变化'", 1),
        ("'追踪政策变化对行业的影响'", 1)
    ])
    
    # Slide 12: 场景选择引导
    slide = add_bullet_slide("思考与选择", [
        "🤔 你对哪个场景最感兴趣？",
        "",
        "💡 课堂练习：选择一个场景进行提问练习",
        "",
        "📝 期末作品：思考你的作品可以做什么方向？",
        "",
        "⏰ 提示：第 6-8 周将深入实践这些场景"
    ])
    
    # ========== 第三部分：高阶提问训练 ==========
    
    # Slide 13: 提问进阶金字塔
    slide = add_bullet_slide("提问进阶金字塔", [
        "🔺 第 5 层：创造性问题（'如果...会怎样'）",
        "",
        "🔺 第 4 层：批判性问题（'这个结论可靠吗'）",
        "",
        "🔺 第 3 层：分析性问题（'为什么/如何'）",
        "",
        "🔺 第 2 层：解释性问题（'什么意思'）",
        "",
        "🔺 第 1 层：事实性问题（'是什么'）"
    ])
    
    # Slide 14: 第 1 层
    slide = add_bullet_slide("第 1 层：事实性问题", [
        "📌 定义：询问基本信息",
        "",
        "✅ 示例：",
        ("'什么是数字人文？'", 1),
        ("'某公司 2024 年营收是多少？'", 1),
        "",
        "⭐ 特点：有明确答案，AI 擅长",
        "",
        "⚠️ 局限：停留在表面"
    ])
    
    # Slide 15: 第 3 层
    slide = add_bullet_slide("第 3 层：分析性问题", [
        "📌 定义：询问原因、机制、过程",
        "",
        "✅ 示例：",
        ("'数字人文如何改变传统文学研究方法？'", 1),
        ("'为什么某品牌在年轻人中更受欢迎？'", 1),
        "",
        "⭐ 特点：需要推理和解释",
        "",
        "💡 价值：深入理解"
    ])
    
    # Slide 16: 第 4 层
    slide = add_bullet_slide("第 4 层：批判性问题", [
        "📌 定义：质疑、验证、评估",
        "",
        "✅ 示例：",
        ("'AI 分析的文本情感结果，如何验证其准确性？'", 1),
        ("'这个结论的样本量是否足够？'", 1),
        "",
        "⭐ 特点：培养批判思维",
        "",
        "💡 价值：避免盲信"
    ])
    
    # Slide 17: 第 5 层
    slide = add_bullet_slide("第 5 层：创造性问题", [
        "📌 定义：探索可能性、创新",
        "",
        "✅ 示例：",
        ("'如果用 AI 分析 100 年的报纸，能发现什么历史规律？'", 1),
        ("'如果结合文本挖掘和知识图谱，能做什么新研究？'", 1),
        "",
        "⭐ 特点：开拓边界",
        "",
        "💡 价值：创新发现"
    ])
    
    # Slide 18: 提问进阶示例（商科）
    slide = add_bullet_slide("提问进阶示例（商科）", [
        "📊 话题：企业社会责任（CSR）",
        "",
        "L1: '什么是企业社会责任？'",
        "L2: '企业社会责任报告包含哪些内容？'",
        "L3: '企业社会责任如何影响消费者购买决策？'",
        "L4: 'AI 分析的企业社会责任报告情感倾向，如何验证？'",
        "L5: '如果用 AI 分析 1000 家企业的 CSR 报告，能发现什么规律？'"
    ])
    
    # Slide 19: 课堂练习说明
    slide = add_bullet_slide("课堂练习：提问进阶训练", [
        "📝 任务：每人完成 3 层提问练习",
        "",
        "✅ 要求：",
        ("选择一个商科话题", 1),
        ("分别写出 L1、L3、L4 层问题", 1),
        ("用 AI 实际提问并记录回答", 1),
        "",
        "⏰ 时间：15 分钟",
        "",
        "📄 提交：课堂练习纸"
    ])
    
    # Slide 20: 练习话题参考
    slide = add_bullet_slide("练习话题参考", [
        "📌 话题 1：某品牌社交媒体评论分析",
        "📌 话题 2：某企业年报情感分析",
        "📌 话题 3：电商产品评论挖掘",
        "📌 话题 4：财经新闻趋势分析",
        "📌 话题 5：自选商科话题",
        "",
        "💡 建议选择你感兴趣的期末作品方向"
    ])
    
    # ========== 第四部分：生态养成研讨 ==========
    
    # Slide 21: 人机共生关系
    slide = add_bullet_slide("你与小龙虾是什么关系？", [
        "🔧 工具关系：你用它干活",
        "",
        "🤝 伙伴关系：互相协作",
        "",
        "🌱 共生关系：共同成长",
        "",
        "🎯 目标：建立健康的共生关系"
    ])
    
    # Slide 22: 小组讨论说明
    slide = add_bullet_slide("小组讨论", [
        "👥 分组：4-5 人/组",
        "",
        "⏰ 时间：20 分钟",
        "",
        "📋 形式：每组选择 1-2 个议题讨论",
        "",
        "🎤 输出：每组派代表分享（3 分钟）"
    ])
    
    # Slide 23: 议题 1
    slide = add_bullet_slide("讨论议题 1：依赖边界", [
        "❓ 核心问题：什么时候应该关掉小龙虾自己思考？",
        "",
        "💡 引导方向：",
        ("考试/测验时", 1),
        ("需要深度思考时", 1),
        ("培养核心能力时", 1),
        ("涉及伦理判断时", 1),
        "",
        "📝 记录要点：小组共识"
    ])
    
    # Slide 24: 议题 2
    slide = add_bullet_slide("讨论议题 2：信任建立", [
        "❓ 核心问题：如何判断 AI 输出是否可信？",
        "",
        "💡 引导方向：",
        ("交叉验证方法", 1),
        ("来源核实", 1),
        ("逻辑检查", 1),
        ("专家意见对比", 1),
        "",
        "📝 记录要点：验证清单"
    ])
    
    # Slide 25: 议题 3
    slide = add_bullet_slide("讨论议题 3：成长记录", [
        "❓ 核心问题：如何记录你和小龙虾的共同成长？",
        "",
        "💡 引导方向：",
        ("互动日志", 1),
        ("能力提升追踪", 1),
        ("作品迭代记录", 1),
        ("反思总结", 1),
        "",
        "📝 记录要点：成长档案设计"
    ])
    
    # Slide 26: 议题 4
    slide = add_bullet_slide("讨论议题 4：伦理困境", [
        "❓ 核心问题：遇到 AI 建议与学术规范冲突怎么办？",
        "",
        "💡 引导方向：",
        ("学术诚信红线", 1),
        ("合理使用边界", 1),
        ("标注规范", 1),
        ("责任归属", 1),
        "",
        "📝 记录要点：伦理准则"
    ])
    
    # Slide 27: 小组分享
    slide = add_bullet_slide("小组分享", [
        "🎤 每组 3 分钟",
        "",
        "📋 分享内容：",
        ("选择的议题", 1),
        ("讨论要点", 1),
        ("形成的共识/建议", 1),
        "",
        "💬 其他组可补充"
    ])
    
    # Slide 28: 教师总结
    slide = add_bullet_slide("健康人机关系的核心原则", [
        "✅ 使用但不依赖",
        "",
        "✅ 信任但验证",
        "",
        "✅ 协作但独立",
        "",
        "✅ 成长但负责",
        "",
        "🎯 最终目标：让你变得更强大，而不是离不开我"
    ])
    
    # ========== 第五部分：期末作品预告 ==========
    
    # Slide 29: 11 周课程全景图
    slide = add_bullet_slide("11 周课程全景图", [
        "📅 第 1 周：认知导入 → 数字人文是什么",
        "📅 第 2-4 周：基础技能 → 小龙虾三部曲",
        "📍 第 5 周：生态认知 → 互动与生态养成（本周）",
        "📅 第 6-8 周：实战深化 → 文本/图谱/可视化",
        "📅 第 9 周：中期检查 → 进度诊断",
        "📅 第 10 周：打磨彩排 → 作品优化",
        "📅 第 11 周：汇报展示 → 成果展示"
    ])
    
    # Slide 30: 期末作品要求
    slide = add_bullet_slide("期末作品要求", [
        "📌 主题：数字人文相关研究或创作",
        "",
        "🤖 AI 使用：必须使用小龙虾协助，并标注使用部分",
        "",
        "💡 个人贡献：必须有独立思考与批判性分析",
        "",
        "📊 形式：报告/可视化/创意作品/工具开发（可选）",
        "",
        "🎤 展示：第 11 周课堂汇报（10 分钟/组）"
    ])
    
    # Slide 31: 选题方向建议
    slide = add_bullet_slide("选题方向建议", [
        "📝 文本分析类",
        ("古诗情感分析、年报挖掘等", 1),
        "",
        "🕸️ 知识图谱类",
        ("人物关系、企业网络等", 1),
        "",
        "📊 可视化类",
        ("数据图表、仪表板等", 1),
        "",
        "✍️ 创意写作类",
        ("人机协作创作", 1),
        "",
        "🔧 工具评测类",
        ("AI 工具对比研究", 1)
    ])
    
    # Slide 32: 商科选题参考
    slide = add_bullet_slide("商科选题参考", [
        "📊 会计学：上市公司年报情感分析",
        "📱 市场营销：品牌评论情感演变",
        "🏢 工商管理：浙商企业家精神研究",
        "🛒 电子商务：电商评论影响分析",
        "💰 金融学：财经新闻与股市关联",
        "",
        "💡 更多参考见《商科选题指南》文档"
    ])
    
    # Slide 33: 个人思考时间
    slide = add_bullet_slide("个人思考时间", [
        "🧘 闭上眼睛想 1 分钟",
        "",
        "💭 思考问题：",
        ("你对什么人文/商科话题感兴趣？", 1),
        ("小龙虾能帮你做什么？", 1),
        ("你想做出什么样的作品？", 1),
        "",
        "⏰ 时间：5 分钟静默思考",
        "",
        "📝 记录：写下初步想法"
    ])
    
    # ========== 第六部分：作业布置 ==========
    
    # Slide 34: 本周作业
    slide = add_bullet_slide("本周作业", [
        "📝 作业 1：互动日志×3",
        ("记录 3 次深度对话（含追问过程）", 1),
        ("使用互动日志模板", 1),
        ("截止时间：第 6 周课前", 1),
        "",
        "📝 作业 2：选题草案",
        ("300 字期末作品初步想法", 1),
        ("包含研究问题", 1),
        ("截止时间：第 6 周课前", 1),
        "",
        "📝 作业 3：反思短文",
        ("'我眼中的健康人机关系'", 1),
        ("300 字", 1),
        ("截止时间：第 6 周课前", 1)
    ])
    
    # Slide 35: 下周预告 & 结束页
    slide = add_title_slide(
        "第 6 周：文本挖掘实战",
        "准备好你的小龙虾，一起探索商业文本的奥秘！\n\n谢谢！🦞"
    )
    
    # 保存文件
    prs.save('/home/admin/.openclaw/workspace/week5_materials/02_第 5 周 PPT.pptx')
    print("PPT 生成完成：第 5 周 PPT.pptx")
    return prs

if __name__ == "__main__":
    create_week5_ppt()
