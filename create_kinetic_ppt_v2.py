#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
人教版高中物理必修 2《动能 动能定理》PPT 课件生成器 v2
改进：素养目标 + 美化排版 + 动画效果
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml

def add_animation(shape, anim_type='appear'):
    """为形状添加动画（简化版）"""
    # python-pptx 动画支持有限，这里添加动画标记
    # 实际动画建议在 PowerPoint 中手动添加
    pass

def create_kinetic_energy_ppt_v2():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # 定义颜色（更丰富的配色）
    PRIMARY = RGBColor(41, 128, 185)      # 主蓝色
    PRIMARY_DARK = RGBColor(30, 90, 140)  # 深蓝
    PRIMARY_LIGHT = RGBColor(174, 214, 241)  # 浅蓝
    ACCENT = RGBColor(230, 126, 34)       # 橙色点缀
    SUCCESS = RGBColor(39, 174, 96)       # 绿色
    BG_LIGHT = RGBColor(248, 249, 250)    # 浅灰背景
    TEXT_DARK = RGBColor(44, 62, 80)      # 深灰文字
    TEXT_MEDIUM = RGBColor(100, 110, 120) # 中灰文字
    
    def add_gradient_bg(slide, color1, color2):
        """添加渐变背景（简化为矩形色块）"""
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.8)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color1
        shape.line.fill.background()
        
        shape2 = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, Inches(1.5), prs.slide_width, Inches(0.5)
        )
        shape2.fill.solid()
        shape2.fill.fore_color.rgb = color2
        shape2.line.fill.background()
        shape2.fill.transparency = 0.7
    
    def add_title_slide(prs, title, subtitle):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # 渐变背景
        add_gradient_bg(slide, PRIMARY, PRIMARY_DARK)
        
        # 装饰圆形
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(11.5), Inches(0.2), Inches(1.5), Inches(1.5)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = ACCENT
        circle.fill.transparency = 0.3
        circle.line.fill.background()
        
        # 标题
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(10), Inches(1.2))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        
        # 副标题
        subtitle_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(10), Inches(0.8))
        tf = subtitle_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(22)
        p.font.color.rgb = PRIMARY_LIGHT
        p.font.italic = True
        
        # 图标
        icon_box = slide.shapes.add_textbox(Inches(0.8), Inches(3.5), Inches(10), Inches(3))
        tf = icon_box.text_frame
        p = tf.paragraphs[0]
        p.text = "⚡  📐"
        p.font.size = Pt(60)
        p.alignment = PP_ALIGN.CENTER
        
        return slide
    
    def add_content_slide(prs, title, content_items, icon="📌"):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # 顶部色块
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.3)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = PRIMARY
        shape.line.fill.background()
        
        # 标题（带图标）
        title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(12), Inches(0.7))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"{icon} {title}"
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        
        # 内容区域背景
        bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.6), Inches(12.3), Inches(5.5)
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = BG_LIGHT
        bg.line.color.rgb = PRIMARY_LIGHT
        bg.line.width = Pt(2)
        
        # 内容
        y_pos = 1.8
        for item in content_items:
            if isinstance(item, dict):
                # 带标题的内容块
                title_shape = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(y_pos), Inches(3.5), Inches(0.5)
                )
                title_shape.fill.solid()
                title_shape.fill.fore_color.rgb = PRIMARY_LIGHT
                title_shape.line.fill.background()
                
                tf = title_shape.text_frame
                p = tf.paragraphs[0]
                p.text = item.get('title', '')
                p.font.size = Pt(18)
                p.font.bold = True
                p.font.color.rgb = PRIMARY_DARK
                
                y_pos += 0.55
                
                for subitem in item.get('items', []):
                    box = slide.shapes.add_textbox(Inches(1.0), Inches(y_pos), Inches(11.5), Inches(0.45))
                    tf = box.text_frame
                    p = tf.paragraphs[0]
                    p.text = f"• {subitem}"
                    p.font.size = Pt(18)
                    p.font.color.rgb = TEXT_MEDIUM
                    y_pos += 0.45
                y_pos += 0.25
            else:
                box = slide.shapes.add_textbox(Inches(0.8), Inches(y_pos), Inches(11.8), Inches(0.5))
                tf = box.text_frame
                p = tf.paragraphs[0]
                p.text = item
                p.font.size = Pt(20)
                p.font.color.rgb = TEXT_DARK
                y_pos += 0.55
        
        return slide
    
    def add_core_competency_slide(prs):
        """核心素养目标页"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # 顶部色块
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.3)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = PRIMARY
        shape.line.fill.background()
        
        # 标题
        title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(12), Inches(0.7))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = "🎯 核心素养目标"
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        
        # 四个素养目标（2x2 网格）
        competencies = [
            ("物理观念", "• 理解动能的概念及其标量性\n• 掌握动能定理的物理意义\n• 认识功与能转化的关系", SUCCESS),
            ("科学思维", "• 通过理论推导建立动能表达式\n• 运用演绎推理探究动能定理\n• 培养模型建构和科学推理能力", ACCENT),
            ("科学探究", "• 经历动能定理的探究过程\n• 学习用数学方法解决物理问题\n• 培养分析论证能力", PRIMARY),
            ("科学态度与责任", "• 感受物理规律的简洁美\n• 培养严谨的科学态度\n• 认识动能定理的实际应用价值", RGBColor(155, 89, 182))
        ]
        
        positions = [
            (Inches(0.6), Inches(1.7), Inches(5.8), Inches(2.5)),
            (Inches(6.9), Inches(1.7), Inches(5.8), Inches(2.5)),
            (Inches(0.6), Inches(4.5), Inches(5.8), Inches(2.5)),
            (Inches(6.9), Inches(4.5), Inches(5.8), Inches(2.5))
        ]
        
        icons = ["⚡", "🧠", "🔬", "💡"]
        
        for i, (comp, desc, color) in enumerate(competencies):
            pos = positions[i]
            
            # 卡片背景
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, pos[0], pos[1], pos[2], pos[3]
            )
            card.fill.solid()
            card.fill.fore_color.rgb = BG_LIGHT
            card.line.color.rgb = color
            card.line.width = Pt(3)
            
            # 图标
            icon_box = slide.shapes.add_textbox(pos[0]+Inches(0.3), pos[1]+Inches(0.2), Inches(1), Inches(1))
            tf = icon_box.text_frame
            p = tf.paragraphs[0]
            p.text = icons[i]
            p.font.size = Pt(40)
            
            # 标题
            title_box = slide.shapes.add_textbox(pos[0]+Inches(1.3), pos[1]+Inches(0.3), Inches(4), Inches(0.6))
            tf = title_box.text_frame
            p = tf.paragraphs[0]
            p.text = comp
            p.font.size = Pt(22)
            p.font.bold = True
            p.font.color.rgb = color
            
            # 描述
            desc_box = slide.shapes.add_textbox(pos[0]+Inches(0.5), pos[1]+Inches(1.1), Inches(5), Inches(1.3))
            tf = desc_box.text_frame
            tf.word_wrap = True
            for j, line in enumerate(desc.split('\n')):
                if j == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = line
                p.font.size = Pt(16)
                p.font.color.rgb = TEXT_MEDIUM
                p.space_after = Pt(6)
        
        return slide
    
    def add_formula_slide(prs, title, formula, explanation, color=PRIMARY):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # 顶部色块
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.3)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = PRIMARY
        shape.line.fill.background()
        
        # 标题
        title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(12), Inches(0.7))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        
        # 公式框（带阴影效果）
        formula_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(1.6), Inches(10.3), Inches(2.2)
        )
        formula_box.fill.solid()
        formula_box.fill.fore_color.rgb = RGBColor(255, 255, 255)
        formula_box.line.color.rgb = color
        formula_box.line.width = Pt(4)
        
        # 添加装饰
        decor = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(11.2), Inches(1.3), Inches(1), Inches(1)
        )
        decor.fill.solid()
        decor.fill.fore_color.rgb = color
        decor.fill.transparency = 0.5
        decor.line.fill.background()
        
        # 公式内容
        tf = formula_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = formula
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = color
        p.alignment = PP_ALIGN.CENTER
        
        # 说明
        y_pos = 4.1
        for line in explanation.split('\n'):
            box = slide.shapes.add_textbox(Inches(0.8), Inches(y_pos), Inches(11.8), Inches(0.5))
            tf = box.text_frame
            p = tf.paragraphs[0]
            p.text = line
            p.font.size = Pt(19)
            p.font.color.rgb = TEXT_MEDIUM
            y_pos += 0.5
        
        return slide
    
    def add_animation_sequence_slide(prs, title, steps):
        """创建带动画效果的演示序列"""
        slides = []
        for i, step in enumerate(steps):
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            
            # 顶部色块
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.3)
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = PRIMARY
            shape.line.fill.background()
            
            # 标题（带进度）
            title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(10), Inches(0.7))
            tf = title_box.text_frame
            p = tf.paragraphs[0]
            p.text = f"🎬 {title}"
            p.font.size = Pt(28)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            
            # 进度指示器
            progress_bg = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(10.5), Inches(0.4), Inches(2.3), Inches(0.5)
            )
            progress_bg.fill.solid()
            progress_bg.fill.fore_color.rgb = RGBColor(255, 255, 255)
            progress_bg.fill.transparency = 0.3
            progress_bg.line.fill.background()
            
            progress_text = slide.shapes.add_textbox(Inches(10.6), Inches(0.45), Inches(2), Inches(0.4))
            tf = progress_text.text_frame
            p = tf.paragraphs[0]
            p.text = f"步骤 {i+1}/{len(steps)}"
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            
            # 步骤进度条
            for j in range(len(steps)):
                seg = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE, 
                    Inches(0.6 + j * 0.65), Inches(1.15), Inches(0.6), Inches(0.15)
                )
                seg.fill.solid()
                if j <= i:
                    seg.fill.fore_color.rgb = ACCENT
                else:
                    seg.fill.fore_color.rgb = PRIMARY_LIGHT
                seg.line.fill.background()
            
            # 内容区域
            content_bg = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.5), Inches(12.1), Inches(5.5)
            )
            content_bg.fill.solid()
            content_bg.fill.fore_color.rgb = BG_LIGHT
            content_bg.line.color.rgb = PRIMARY_LIGHT
            content_bg.line.width = Pt(2)
            
            # 步骤内容
            y_pos = 1.8
            for line in step.split('\n'):
                if line.strip():
                    # 高亮关键行
                    is_highlight = any(kw in line for kw in ['关键', '→', '【步骤'])
                    
                    if is_highlight:
                        highlight = slide.shapes.add_shape(
                            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(y_pos-0.1), Inches(11.7), Inches(0.5)
                        )
                        highlight.fill.solid()
                        highlight.fill.fore_color.rgb = ACCENT
                        highlight.fill.transparency = 0.2
                        highlight.line.fill.background()
                    
                    box = slide.shapes.add_textbox(Inches(0.9), Inches(y_pos), Inches(11.5), Inches(0.5))
                    tf = box.text_frame
                    p = tf.paragraphs[0]
                    p.text = line
                    p.font.size = Pt(20)
                    if is_highlight:
                        p.font.bold = True
                        p.font.color.rgb = RGBColor(200, 80, 0)
                    else:
                        p.font.color.rgb = TEXT_MEDIUM
                    y_pos += 0.55
            
            slides.append(slide)
        
        return slides
    
    # ========== 创建幻灯片 ==========
    
    # 第 1 页：封面
    add_title_slide(prs, "动能 动能定理", "人教版高中物理必修 2 · 第七章 机械能守恒定律")
    
    # 第 2 页：核心素养目标
    add_core_competency_slide(prs)
    
    # 第 3 页：情境导入
    add_content_slide(prs, "情境导入", [
        "思考以下问题：",
        "• 为什么高速行驶的汽车刹车距离更长？",
        "• 为什么子弹虽然质量小，但杀伤力很大？",
        "• 物体的动能与哪些因素有关？",
        "",
        "【生活实例】",
        "• 流星撞击地球释放巨大能量  ☄️",
        "• 风力发电机利用空气动能发电  🌬️",
        "• 打桩机利用重锤的动能做功  🏗️"
    ], "💡")
    
    # 第 4 页：动能的概念
    add_content_slide(prs, "动能的概念", [
        {"title": "定义", "items": ["物体由于运动而具有的能量"]},
        {"title": "表达式", "items": ["Eₖ = ½mv²", "m — 质量 (kg)", "v — 瞬时速度 (m/s)", "Eₖ — 动能 (J)"]},
        {"title": "单位", "items": ["焦耳 (J)", "1 J = 1 kg·m²/s² = 1 N·m"]},
        {"title": "特点", "items": ["标量，只有大小，没有方向", "总是正值 (v²≥0)", "具有相对性，与参考系有关"]}
    ], "📖")
    
    # 第 5 页：动能表达式推导
    add_formula_slide(prs, "动能表达式的推导", 
        "Eₖ = ½mv²",
        "推导思路：从功与能的关系出发\n\n设质量为 m 的物体，初速度为 v₀，在恒力 F 作用下\n经过位移 l，末速度为 v",
        SUCCESS
    )
    
    # 第 6-9 页：动能定理推导演示动画序列
    anim_steps = [
        "【步骤 1】物理模型\n\n• 光滑水平面上，质量为 m 的物体\n• 初速度 v₀，受恒力 F 作用\n• 经过位移 l，末速度为 v\n\n→ 分析力做功与速度变化的关系",
        
        "【步骤 2】牛顿第二定律\n\n• 物体受力：F（合力）\n• 加速度：a = F/m\n• 物体做匀加速直线运动\n\n→ 建立力与加速度的关系",
        
        "【步骤 3】运动学公式\n\n• 匀变速直线运动：v² - v₀² = 2al\n• 变形得：l = (v² - v₀²)/(2a)\n\n→ 建立位移与速度的关系",
        
        "【关键步骤 4】功的计算\n\n• 合力做功：W = F·l\n• 代入 l：W = F·(v² - v₀²)/(2a)\n• 代入 a = F/m：W = ½mv² - ½mv₀²\n\n→ 得到功与动能变化的关系"
    ]
    
    add_animation_sequence_slide(prs, "动能定理推导", anim_steps)
    
    # 第 10 页：动能定理
    add_formula_slide(prs, "动能定理", 
        "W = Eₖ₂ - Eₖ₁ = ΔEₖ",
        "内容：合力对物体所做的功，等于物体动能的变化量\n\n物理意义：\n• 揭示了功与能的转化关系\n• 功是能量转化的量度\n• 动能的变化由合力做功决定",
        ACCENT
    )
    
    # 第 11 页：动能定理的理解
    add_content_slide(prs, "动能定理的理解", [
        {"title": "适用范围", "items": ["恒力做功、变力做功", "直线运动、曲线运动", "单个物体、物体系"]},
        {"title": "W 的含义", "items": ["合力做的功（或各力做功代数和）", "W = W₁ + W₂ + W₃ + ...", "包括重力、弹力、摩擦力等所有力"]},
        {"title": "解题优势", "items": ["不涉及加速度和时间，简化计算", "特别适用于变力做功和曲线运动", "解决力学问题的重要工具"]}
    ], "💭")
    
    # 第 12 页：应用示例 1
    add_content_slide(prs, "应用示例 1：刹车距离", [
        {"title": "题目", "items": ["质量 m = 1500 kg 的汽车，v = 20 m/s", "刹车阻力 f = 6000 N", "求：刹车距离"]},
        {"title": "解析", "items": ["由动能定理：W = ΔEₖ", "-f·s = 0 - ½mv²", "s = mv²/(2f) = 1500×20²/(2×6000) = 50 m"]},
        {"title": "思考", "items": ["速度加倍，刹车距离变为多少？（4 倍！）", "这就是为什么不能超速行驶 ⚠️"]}
    ], "✏️")
    
    # 第 13 页：应用示例 2
    add_content_slide(prs, "应用示例 2：自由落体", [
        {"title": "题目", "items": ["质量为 m 的物体从高度 h 处自由下落", "求：落地时的速度（不计空气阻力）"]},
        {"title": "解析", "items": ["由动能定理：W = ΔEₖ", "mgh = ½mv² - 0", "v = √(2gh)"]},
        {"title": "说明", "items": ["与运动学公式结果一致", "动能定理更简洁，不涉及时间"]}
    ], "✏️")
    
    # 第 14 页：课堂练习
    add_content_slide(prs, "课堂练习", [
        {"title": "基础题", "items": ["质量 2 kg 的物体，速度从 3 m/s 增加到 5 m/s", "求动能的变化量", "【答案】ΔEₖ = ½×2×(5²-3²) = 16 J"]},
        {"title": "提升题", "items": ["物体从斜面顶端滑下，h = 5 m，l = 10 m", "摩擦因数μ = 0.2，求到达底端的速度", "【提示】mgh - μmgcosθ·l = ½mv²"]},
        {"title": "拓展题", "items": ["思考：为什么过山车要从高处开始？", "用动能定理解释"]}
    ], "📝")
    
    # 第 15 页：易错点警示
    add_content_slide(prs, "易错点警示", [
        {"title": "常见错误", "items": ["忘记动能是标量，误认为有方向", "混淆动能和动量 (mv)", "计算合力功时漏掉某个力", "W 是合力功，不是某个力的功"]},
        {"title": "注意事项", "items": ["明确研究对象和过程", "正确分析受力，计算各力做功", "注意初末状态的动能", "单位统一用国际单位制"]}
    ], "⚠️")
    
    # 第 16 页：知识总结
    add_content_slide(prs, "知识总结", [
        {"title": "核心公式", "items": ["动能：Eₖ = ½mv²", "动能定理：W = ΔEₖ = Eₖ₂ - Eₖ₁"]},
        {"title": "物理思想", "items": ["能量观：功是能量转化的量度", "状态观：动能是状态量，功是过程量", "转化观：不同形式能量可以相互转化"]},
        {"title": "解题步骤", "items": ["1. 确定研究对象和过程", "2. 分析受力，计算合力做功", "3. 确定初末动能", "4. 列动能定理方程求解"]}
    ], "📊")
    
    # 第 17 页：课后作业
    add_content_slide(prs, "课后作业", [
        "【必做题】教材 P73 练习 1、2、3",
        "",
        "【选做题】",
        "• 查阅资料：动能定理在生活中的应用",
        "• 思考题：如果考虑空气阻力，自由落体的",
        "  落地速度如何计算？",
        "",
        "【预习】",
        "• 预习下一节：机械能守恒定律",
        "• 思考：什么情况下机械能守恒？"
    ], "📚")
    
    # 第 18 页：结束页
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 背景
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, Inches(2), prs.slide_width, Inches(5.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = PRIMARY
    shape.line.fill.background()
    
    # 装饰
    circle1 = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(10), Inches(2.5), Inches(2), Inches(2)
    )
    circle1.fill.solid()
    circle1.fill.fore_color.rgb = ACCENT
    circle1.fill.transparency = 0.5
    circle1.line.fill.background()
    
    circle2 = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(11.5), Inches(5), Inches(1.5), Inches(1.5)
    )
    circle2.fill.solid()
    circle2.fill.fore_color.rgb = SUCCESS
    circle2.fill.transparency = 0.5
    circle2.line.fill.background()
    
    # 文字
    box = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(12.3), Inches(2))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = "谢谢观看！"
    p.font.size = Pt(52)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    p2 = tf.add_paragraph()
    p2.text = "⚡ 勤学善思 · 格物致知 ⚡"
    p2.font.size = Pt(26)
    p2.font.color.rgb = PRIMARY_LIGHT
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(20)
    
    # 保存
    output_path = "/home/admin/.openclaw/workspace/动能_动能定理_v2.pptx"
    prs.save(output_path)
    
    return output_path

if __name__ == "__main__":
    path = create_kinetic_energy_ppt_v2()
    print(f"PPT v2 已生成：{path}")
