#!/usr/bin/env python3
"""生成"运营商数字员工"科普PPT"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ===== 配色方案 =====
DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)       # 深蓝黑背景
ACCENT_BLUE = RGBColor(0x00, 0x7B, 0xFF)    # 亮蓝
ACCENT_CYAN = RGBColor(0x00, 0xD4, 0xFF)    # 青色
ACCENT_ORANGE = RGBColor(0xFF, 0x8C, 0x00)  # 橙色
ACCENT_GREEN = RGBColor(0x00, 0xC8, 0x96)   # 绿色
ACCENT_RED = RGBColor(0xFF, 0x4D, 0x4D)     # 红色
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
MID_GRAY = RGBColor(0x99, 0x99, 0x99)
CARD_BG = RGBColor(0x25, 0x25, 0x40)        # 卡片背景
SUBTITLE_GRAY = RGBColor(0xAA, 0xAA, 0xBB)

OUTPUT_DIR = os.path.dirname(os.path.abspath(__file__))


def add_bg(slide, color=DARK_BG):
    """添加纯色背景"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color, shape_type=MSO_SHAPE.RECTANGLE):
    """添加形状"""
    shape = slide.shapes.add_shape(shape_type, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="微软雅黑"):
    """添加文本框"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_card(slide, left, top, width, height, title, body_lines, title_color=ACCENT_BLUE, icon=""):
    """添加卡片"""
    # 卡片背景
    card = add_shape(slide, left, top, width, height, CARD_BG)
    card.shadow.inherit = False

    # 标题
    title_y = top + 0.2
    add_text_box(slide, left + 0.3, title_y, width - 0.6, 0.5, icon + title, font_size=16, color=title_color, bold=True)

    # 正文
    body_y = title_y + 0.55
    txBox = slide.shapes.add_textbox(Inches(left + 0.3), Inches(body_y), Inches(width - 0.6), Inches(height - 0.8))
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, line in enumerate(body_lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(13)
        p.font.color.rgb = LIGHT_GRAY
        p.font.name = "微软雅黑"
        p.space_after = Pt(6)

    return card


def add_divider_line(slide, left, top, width, color=ACCENT_BLUE, thickness=2):
    """添加分隔线"""
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(thickness / 72.0)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()


# ============================================================
# 第1页：封面
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白版式
add_bg(slide)

# 顶部装饰线
add_shape(slide, 0, 0, 13.333, 0.08, ACCENT_BLUE)

# 主标题
add_text_box(slide, 1.5, 1.8, 10.3, 1.5,
             "运营商同时上线\n「数字员工」",
             font_size=52, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# 副标题
add_text_box(slide, 1.5, 3.8, 10.3, 0.8,
             "中国移动 8 万数智员工 · 中国电信 50 类数字员工",
             font_size=24, color=ACCENT_CYAN, bold=False, alignment=PP_ALIGN.CENTER)

# 分隔线
add_divider_line(slide, 5, 4.8, 3.333, ACCENT_CYAN, 3)

# 说明文字
add_text_box(slide, 1.5, 5.2, 10.3, 1.0,
             "AI 应用层的争夺，已经正式开始",
             font_size=20, color=SUBTITLE_GRAY, bold=False, alignment=PP_ALIGN.CENTER)

# 底部标签
add_text_box(slide, 1.5, 6.3, 10.3, 0.5,
             "🦞 诸葛虾 · 科普系列    |    2026年5月",
             font_size=14, color=MID_GRAY, bold=False, alignment=PP_ALIGN.CENTER)


# ============================================================
# 第2页：什么是「数字员工」？
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, 1.0, 0.5, 11.3, 0.8, "01  什么是「数字员工」？", font_size=32, color=ACCENT_BLUE, bold=True)
add_divider_line(slide, 1.0, 1.2, 3, ACCENT_BLUE, 3)

# 定义卡片
add_card(slide, 1.0, 1.7, 11.3, 1.5,
         "简单来说",
         ["数字员工 = AI 驱动的虚拟员工，能帮你干活、做决策、处理流程。"],
         title_color=ACCENT_CYAN, icon="💡 ")

# 三个示例卡片
add_card(slide, 1.0, 3.6, 3.6, 2.8,
         "📞 客服数字员工",
         ["7×24 小时在线", "自动回答客户问题", "识别情绪并转人工"],
         title_color=ACCENT_ORANGE)

add_card(slide, 5.0, 3.6, 3.6, 2.8,
         "📊 财务数字员工",
         ["自动记账对账", "发票识别审核", "生成财务报表"],
         title_color=ACCENT_GREEN)

add_card(slide, 9.0, 3.6, 3.6, 2.8,
         "🔧 运维数字员工",
         ["监控网络故障", "自动排查问题", "预测设备风险"],
         title_color=ACCENT_CYAN)

# 底部提示
add_text_box(slide, 1.0, 6.8, 11.3, 0.5,
             "不是科幻电影——它们已经在我们身边了",
             font_size=14, color=MID_GRAY, bold=False, alignment=PP_ALIGN.CENTER)


# ============================================================
# 第3页：新闻速览
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, 1.0, 0.5, 11.3, 0.8, "02  新闻速览：两大运营商同时出手", font_size=32, color=ACCENT_BLUE, bold=True)
add_divider_line(slide, 1.0, 1.2, 3.5, ACCENT_BLUE, 3)

# 中国移动卡片
add_card(slide, 1.0, 1.7, 5.3, 4.5,
         "🔵 中国移动",
         ["📍 数字中国峰会宣布",
          "👥 上线 8 万数智员工",
          "🏢 全部用于内部：网络鉴伪、研发设计、营销服务、综合管理",
          "🎯 目标：降本增效，刀刃向内"],
         title_color=ACCENT_CYAN)

# 中国电信卡片
add_card(slide, 7.0, 1.7, 5.3, 4.5,
         "🔴 中国电信",
         ["📍 同期发布",
          "📦 首期 8 类数字员工套餐",
          "🏢 面向中小微企业：行政、财务、客服、销售等",
          "🎯 年底扩展至 50 类，对外销售"],
         title_color=ACCENT_RED)

# 底部关键信息
add_text_box(slide, 1.0, 6.5, 11.3, 0.8,
             "⚡ 两大运营商同时出手，这不是巧合——AI 应用层的争夺，已经正式开始。",
             font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)


# ============================================================
# 第4页：为什么是现在？
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, 1.0, 0.5, 11.3, 0.8, "03  为什么是现在？", font_size=32, color=ACCENT_BLUE, bold=True)
add_divider_line(slide, 1.0, 1.2, 2.5, ACCENT_BLUE, 3)

# 数据卡片
add_text_box(slide, 1.0, 1.7, 11.3, 0.6, "📊 2026年第一季度数据", font_size=22, color=WHITE, bold=True)

# 移动数据
add_card(slide, 1.0, 2.5, 5.3, 3.0,
         "中国移动",
         ["营收 2665 亿元 ↑ 1.0%",
          "净利润 293 亿元 ↓ 4.2%",
          "",
          "💰 营收在涨，利润在跌"],
         title_color=ACCENT_CYAN)

# 电信数据
add_card(slide, 7.0, 2.5, 5.3, 3.0,
         "中国电信",
         ["营收 1314 亿元 ↑ 2.3%",
          "净利润 73.5 亿元 ↓ 17.1%",
          "",
          "💰 同样的故事：成本压力 > 收入增长"],
         title_color=ACCENT_RED)

# 结论
add_card(slide, 1.0, 5.9, 11.3, 1.2,
         "🔍 这意味着什么？",
         ["成本端的压力，已经压过了收入端的增长。推数字员工，不完全是进攻，更多是防守。"],
         title_color=ACCENT_ORANGE, icon="⚠️ ")


# ============================================================
# 第5页：两种打法，两种逻辑
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, 1.0, 0.5, 11.3, 0.8, "04  两种打法，两种逻辑", font_size=32, color=ACCENT_BLUE, bold=True)
add_divider_line(slide, 1.0, 1.2, 3, ACCENT_BLUE, 3)

# 向左箭头
arrow = slide.shapes.add_shape(MSO_SHAPE.LEFT_ARROW, Inches(5.8), Inches(2.5), Inches(1.7), Inches(0.8))
arrow.fill.solid()
arrow.fill.fore_color.rgb = ACCENT_ORANGE
arrow.line.fill.background()

# 移动 - 向内
add_card(slide, 1.0, 2.0, 4.5, 4.5,
         "🏠 中国移动：向内",
         ["练内功，刀刃向内",
          "8 万数字员工全部自用",
          "当成管理工具，不是商业产品",
          "",
          "✅ 优势：场景真实、数据充足",
          "❌ 风险：只节流、不开源"],
         title_color=ACCENT_CYAN)

# 电信 - 向外
add_card(slide, 7.8, 2.0, 4.5, 4.5,
         "🌐 中国电信：向外",
         ["面向中小微企业卖套餐",
          "「数字员工 + Token + 连接」",
          "低成本、快速应用",
          "",
          "✅ 优势：直接创收、抢占市场",
          "❌ 风险：场景深度不够"],
         title_color=ACCENT_RED)


# ============================================================
# 第6页：Token —— 新的计费单位
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, 1.0, 0.5, 11.3, 0.8, "05  Token：AI 时代的新计量单位", font_size=32, color=ACCENT_BLUE, bold=True)
add_divider_line(slide, 1.0, 1.2, 3.5, ACCENT_BLUE, 3)

# 类比卡片
add_card(slide, 1.0, 1.7, 5.3, 2.5,
         "📱 过去：卖流量",
         ["你上网越多 → 流量消耗越多 → 运营商赚得越多",
          "",
          "流量 = 互联网时代的计量单位"],
         title_color=ACCENT_GREEN)

add_card(slide, 7.0, 1.7, 5.3, 2.5,
         "🤖 现在：卖 Token",
         ["你用的数字员工越多 → Token 消耗越多 → 收入越多",
          "",
          "Token = AI 时代的计量单位"],
         title_color=ACCENT_ORANGE)

# 核心逻辑
add_text_box(slide, 1.0, 4.5, 11.3, 0.6, "💡 核心逻辑", font_size=22, color=WHITE, bold=True)

add_card(slide, 1.0, 5.2, 11.3, 1.8,
         "每一个数字员工背后，都是持续的 Token 消耗 + 持续的连接需求",
         ["这个市场的天花板，远比想象中高。",
          "数字员工用得越多 → Token 消耗越多 → 运营商收入越多 → 形成正向循环"],
         title_color=ACCENT_CYAN)


# ============================================================
# 第7页：挑战与难点
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, 1.0, 0.5, 11.3, 0.8, "06  挑战与难点", font_size=32, color=ACCENT_BLUE, bold=True)
add_divider_line(slide, 1.0, 1.2, 2, ACCENT_BLUE, 3)

# 挑战1
add_card(slide, 1.0, 1.7, 5.3, 2.2,
         "移动：从内部到外部",
         ["8 万数智员工是极好的测试场景",
          "但从内部工具到商业产品，这一步需要时间",
          "时间窗口还有多大？竞争对手不会等。"],
         title_color=ACCENT_RED, icon="⏰ ")

# 挑战2
add_card(slide, 7.0, 1.7, 5.3, 2.2,
         "电信：场景深度不够",
         ["中小微企业要的不是概念，是解决问题",
          "数字员工能不能真正帮客户省下真金白银？",
          "套餐卖得出去 ≠ 客户留得住"],
         title_color=ACCENT_RED, icon="🎯 ")

# 挑战3
add_card(slide, 1.0, 4.3, 5.3, 2.2,
         "共同难题：人才",
         ["既懂 AI 又懂行业的人才稀缺",
          "运营商传统上是做管道的，做应用是另一回事"],
         title_color=ACCENT_ORANGE, icon="👥 ")

# 挑战4
add_card(slide, 7.0, 4.3, 5.3, 2.2,
         "信任问题",
         ["企业敢把核心业务交给 AI 吗？",
          "数据安全、隐私保护是关键",
          "运营商的品牌信任度是优势"],
         title_color=ACCENT_ORANGE, icon="🔒 ")

# 底部
add_text_box(slide, 1.0, 6.8, 11.3, 0.5,
             "方向对了，不代表能拿下市场。这中间的距离，不是一两步能跨过去的。",
             font_size=15, color=SUBTITLE_GRAY, bold=False, alignment=PP_ALIGN.CENTER)


# ============================================================
# 第8页：对我们意味着什么？
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, 1.0, 0.5, 11.3, 0.8, "07  对我们意味着什么？", font_size=32, color=ACCENT_BLUE, bold=True)
add_divider_line(slide, 1.0, 1.2, 3, ACCENT_BLUE, 3)

# 三个影响卡片
add_card(slide, 1.0, 1.7, 3.6, 4.5,
         "👔 打工人",
         ["部分重复性工作",
          "会被 AI 替代",
          "",
          "但 AI 更多是助手，",
          "不是替代者。",
          "学会和 AI 协作",
          "才是未来竞争力。"],
         title_color=ACCENT_ORANGE)

add_card(slide, 4.9, 1.7, 3.6, 4.5,
         "🏪 小老板",
         ["中小微企业",
          "可以用更低的成本",
          "获得 AI 员工",
          "",
          "客服、财务、销售",
          "都能找到数字员工",
          "帮你干活。"],
         title_color=ACCENT_GREEN)

add_card(slide, 8.8, 1.7, 3.6, 4.5,
         "🎓 学生",
         ["AI 应用层",
          "是未来最大的就业市场",
          "",
          "懂 AI + 懂行业",
          "= 最吃香的人才",
          "",
          "现在入局，正当时。"],
         title_color=ACCENT_CYAN)

# 底部
add_text_box(slide, 1.0, 6.5, 11.3, 0.8,
             "AI 时代，每个人都需要思考：我如何与 AI 协作？",
             font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)


# ============================================================
# 第9页：未来展望
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, 1.0, 0.5, 11.3, 0.8, "08  未来展望", font_size=32, color=ACCENT_BLUE, bold=True)
add_divider_line(slide, 1.0, 1.2, 2, ACCENT_BLUE, 3)

# 时间线
timeline_y = 2.0

# 短期
add_shape(slide, 1.0, timeline_y, 3.5, 0.15, ACCENT_ORANGE)
add_text_box(slide, 1.0, timeline_y + 0.3, 3.5, 0.5, "短期（1-2年）", font_size=18, color=ACCENT_ORANGE, bold=True)
add_card(slide, 1.0, timeline_y + 0.8, 3.5, 2.2,
         "",
         ["数字员工以「降本」为主",
          "电信套餐会经历客户流失",
          "场景深度决定去留"],
         title_color=ACCENT_ORANGE)

# 中期
add_shape(slide, 4.9, timeline_y, 3.5, 0.15, ACCENT_CYAN)
add_text_box(slide, 4.9, timeline_y + 0.3, 3.5, 0.5, "中期（3-5年）", font_size=18, color=ACCENT_CYAN, bold=True)
add_card(slide, 4.9, timeline_y + 0.8, 3.5, 2.2,
         "",
         ["移动内部验证成熟后",
          "可能产品化推向市场",
          "数据积累成关键优势"],
         title_color=ACCENT_CYAN)

# 长期
add_shape(slide, 8.8, timeline_y, 3.5, 0.15, ACCENT_GREEN)
add_text_box(slide, 8.8, timeline_y + 0.3, 3.5, 0.5, "长期（5年+）", font_size=18, color=ACCENT_GREEN, bold=True)
add_card(slide, 8.8, timeline_y + 0.8, 3.5, 2.2,
         "",
         ["Token 计费成为",
          "运营商新增长曲线",
          "AI 员工成为基础设施"],
         title_color=ACCENT_GREEN)

# 总结
add_card(slide, 1.0, 5.5, 11.3, 1.5,
         "🦞 诸葛虾说",
         ["真正能跑出来的，不是喊得响的，是能把一个场景做到客户离不开的。",
          "数字员工这张牌，两大运营商都押上去了，我们拭目以待。"],
         title_color=ACCENT_CYAN)


# ============================================================
# 第10页：结尾
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_shape(slide, 0, 0, 13.333, 0.08, ACCENT_BLUE)

add_text_box(slide, 1.5, 2.0, 10.3, 1.5,
             "AI 时代，",
             font_size=48, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_text_box(slide, 1.5, 3.0, 10.3, 1.0,
             "每个人都需要学会与 AI 协作。",
             font_size=36, color=ACCENT_CYAN, bold=True, alignment=PP_ALIGN.CENTER)

add_divider_line(slide, 5, 4.3, 3.333, ACCENT_CYAN, 3)

add_text_box(slide, 1.5, 4.8, 10.3, 0.8,
             "数据来源：「运营商那些事」公众号 · 观海",
             font_size=14, color=MID_GRAY, bold=False, alignment=PP_ALIGN.CENTER)

add_text_box(slide, 1.5, 5.5, 10.3, 0.5,
             "🦞 诸葛虾 · 科普系列    |    2026年5月",
             font_size=16, color=SUBTITLE_GRAY, bold=False, alignment=PP_ALIGN.CENTER)

add_text_box(slide, 1.5, 6.3, 10.3, 0.5,
             "感谢观看 ✨",
             font_size=20, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)


# 保存
output_path = os.path.join(OUTPUT_DIR, "运营商数字员工科普.pptx")
prs.save(output_path)
print(f"✅ PPT 已保存: {output_path}")
print(f"📊 共 {len(prs.slides)} 页")
