#!/usr/bin/env python3
"""生成八年级期中考试成绩分析报告的Word和PPT文件"""

from docx import Document
from docx.shared import Pt, Inches, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor as PPTColor
from pptx.enum.text import PP_ALIGN
import os

# ==================== 数据定义（从图片提取） ====================
# 八年级8个班级数据
classes_data = [
    {"id": 1, "students": 380, "score_400": 34, "score_400_pct": 8.95, "score_360": 85, "score_360_pct": 22.37,
     "subjects": {
         "语文": {"avg": 65.4, "excellent": 3.70, "pass": 73.90},
         "数学": {"avg": 59.3, "excellent": 24.50, "pass": 54.50},
         "英语": {"avg": 41.7, "excellent": 3.20, "pass": 18.90},
         "科学": {"avg": 50.3, "excellent": 8.70, "pass": 33.20},
         "社会": {"avg": 58.6, "excellent": 7.40, "pass": 52.60},
     },
     "total": {"avg": 275.3, "excellent": 5.30, "pass": 40.00}},
    
    {"id": 2, "students": 287, "score_400": 50, "score_400_pct": 17.42, "score_360": 88, "score_360_pct": 30.66,
     "subjects": {
         "语文": {"avg": 66.96, "excellent": 5.23, "pass": 78.40},
         "数学": {"avg": 62.7, "excellent": 28.92, "pass": 59.58},
         "英语": {"avg": 47.17, "excellent": 9.41, "pass": 31.36},
         "科学": {"avg": 54.05, "excellent": 11.50, "pass": 43.55},
         "社会": {"avg": 63.22, "excellent": 14.63, "pass": 63.41},
     },
     "total": {"avg": 294.09, "excellent": 9.06, "pass": 52.61}},
    
    {"id": 3, "students": 282, "score_400": 35, "score_400_pct": 12.41, "score_360": 83, "score_360_pct": 29.43,
     "subjects": {
         "语文": {"avg": 63.11, "excellent": 3.55, "pass": 69.86},
         "数学": {"avg": 62.91, "excellent": 27.30, "pass": 63.83},
         "英语": {"avg": 46.03, "excellent": 8.51, "pass": 33.33},
         "科学": {"avg": 57.12, "excellent": 13.48, "pass": 50.35},
         "社会": {"avg": 52.94, "excellent": 4.26, "pass": 53.04},
     },
     "total": {"avg": 282.11, "excellent": 5.32, "pass": 48.23}},
    
    {"id": 4, "students": 355, "score_400": 30, "score_400_pct": 8.45, "score_360": 57, "score_360_pct": 16.06,
     "subjects": {
         "语文": {"avg": 56.03, "excellent": 0.85, "pass": 47.32},
         "数学": {"avg": 52.12, "excellent": 21.13, "pass": 47.32},
         "英语": {"avg": 39.12, "excellent": 1.97, "pass": 16.62},
         "科学": {"avg": 48.10, "excellent": 6.48, "pass": 30.70},
         "社会": {"avg": 57.05, "excellent": 8.17, "pass": 52.39},
     },
     "total": {"avg": 252.42, "excellent": 2.54, "pass": 36.34}},
    
    {"id": 5, "students": 286, "score_400": 28, "score_400_pct": 9.79, "score_360": 52, "score_360_pct": 18.18,
     "subjects": {
         "语文": {"avg": 67.74, "excellent": 6.60, "pass": 81.47},
         "数学": {"avg": 51.81, "excellent": 15.73, "pass": 45.45},
         "英语": {"avg": 44.26, "excellent": 6.64, "pass": 24.83},
         "科学": {"avg": 50.06, "excellent": 9.79, "pass": 33.92},
         "社会": {"avg": 49.54, "excellent": 3.15, "pass": 30.77},
     },
     "total": {"avg": 263.41, "excellent": 3.85, "pass": 35.31}},
    
    {"id": 6, "students": 323, "score_400": 37, "score_400_pct": 14.46, "score_360": 80, "score_360_pct": 24.77,
     "subjects": {
         "语文": {"avg": 63.3, "excellent": 9.60, "pass": 67.80},
         "数学": {"avg": 51.6, "excellent": 19.20, "pass": 43.70},
         "英语": {"avg": 43.4, "excellent": 3.40, "pass": 26.10},
         "科学": {"avg": 51.3, "excellent": 9.60, "pass": 39.30},
         "社会": {"avg": 61, "excellent": 11.50, "pass": 59.80},
     },
     "total": {"avg": 270.5, "excellent": 5.90, "pass": 40.90}},
    
    {"id": 7, "students": 330, "score_400": 39, "score_400_pct": 11.82, "score_360": 77, "score_360_pct": 23.33,
     "subjects": {
         "语文": {"avg": 64.53, "excellent": 3.94, "pass": 72.73},
         "数学": {"avg": 56.11, "excellent": 23.94, "pass": 49.39},
         "英语": {"avg": 46.6, "excellent": 7.27, "pass": 31.21},
         "科学": {"avg": 51.46, "excellent": 12.42, "pass": 36.36},
         "社会": {"avg": 53.44, "excellent": 3.94, "pass": 40.91},
     },
     "total": {"avg": 272.14, "excellent": 5.15, "pass": 41.21}},
    
    {"id": 8, "students": 394, "score_400": 46, "score_400_pct": 11.68, "score_360": 94, "score_360_pct": 23.86,
     "subjects": {
         "语文": {"avg": 66.06, "excellent": 3.35, "pass": 77.06},
         "数学": {"avg": 59.44, "excellent": 22.42, "pass": 57.22},
         "英语": {"avg": 46.2, "excellent": 6.96, "pass": 27.06},
         "科学": {"avg": 56.82, "excellent": 11.34, "pass": 48.97},
         "社会": {"avg": 57.42, "excellent": 5.41, "pass": 49.48},
     },
     "total": {"avg": 285.94, "excellent": 5.93, "pass": 48.71}},
]

OUTPUT_DIR = "/home/admin/.openclaw/workspace"

# 计算整体数据
total_students = sum(c["students"] for c in classes_data)
total_400 = sum(c["score_400"] for c in classes_data)
total_360 = sum(c["score_360"] for c in classes_data)
avg_400_pct = sum(c["score_400_pct"] * c["students"] for c in classes_data) / total_students
avg_360_pct = sum(c["score_360_pct"] * c["students"] for c in classes_data) / total_students

# 计算各科加权平均
subject_avgs = {}
for subj in ["语文", "数学", "英语", "科学", "社会"]:
    weighted_avg = sum(c["subjects"][subj]["avg"] * c["students"] for c in classes_data) / total_students
    excellent = sum(c["subjects"][subj]["excellent"] * c["students"] for c in classes_data) / total_students
    passing = sum(c["subjects"][subj]["pass"] * c["students"] for c in classes_data) / total_students
    subject_avgs[subj] = {"avg": round(weighted_avg, 1), "excellent": round(excellent, 1), "pass": round(passing, 1)}

# 计算总分加权平均
total_avg = sum(c["total"]["avg"] * c["students"] for c in classes_data) / total_students
total_excellent = sum(c["total"]["excellent"] * c["students"] for c in classes_data) / total_students
total_pass = sum(c["total"]["pass"] * c["students"] for c in classes_data) / total_students

overall = {
    "total_students": total_students,
    "score_400": total_400,
    "score_400_pct": round(avg_400_pct, 2),
    "score_360": total_360,
    "score_360_pct": round(avg_360_pct, 2),
    "total": {"avg": round(total_avg, 1), "excellent": round(total_excellent, 2), "pass": round(total_pass, 2)},
    "subjects": subject_avgs,
}

# 按总分平均分排序
ranked_classes = sorted(classes_data, key=lambda x: x["total"]["avg"], reverse=True)


# ==================== Word 文档生成 ====================
def create_word_report():
    doc = Document()
    
    style = doc.styles['Normal']
    font = style.font
    font.name = '微软雅黑'
    font.size = Pt(11)
    
    # 标题
    title = doc.add_heading('八年级 期中考试成绩分析报告', level=0)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    for run in title.runs:
        run.font.color.rgb = RGBColor(0x1a, 0x5c, 0x8a)
    
    doc.add_paragraph('')
    
    # 一、整体概况
    doc.add_heading('一、整体概况', level=1)
    
    table = doc.add_table(rows=4, cols=2, style='Light Shading Accent 1')
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    data = [
        ('班级数', '8个班'),
        ('参考人数', f"{overall['total_students']}人"),
        ('≥400分', f"{overall['score_400']}人（占比{overall['score_400_pct']}%）"),
        ('≥360分', f"{overall['score_360']}人（占比{overall['score_360_pct']}%）"),
    ]
    for i, (key, value) in enumerate(data):
        table.rows[i].cells[0].text = key
        table.rows[i].cells[1].text = value
        for cell in table.rows[i].cells:
            for paragraph in cell.paragraphs:
                paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
    
    doc.add_paragraph('')
    
    # 二、各科成绩详情
    doc.add_heading('二、各科成绩详情（整体）', level=1)
    
    headers = ['科目', '平均分', '优秀率', '及格率']
    table = doc.add_table(rows=6, cols=4, style='Light Shading Accent 1')
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    
    for j, header in enumerate(headers):
        cell = table.rows[0].cells[j]
        cell.text = header
        for paragraph in cell.paragraphs:
            paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
    
    for i, (subject, data) in enumerate(overall['subjects'].items()):
        row = table.rows[i+1]
        row.cells[0].text = subject
        row.cells[1].text = str(data['avg'])
        row.cells[2].text = f"{data['excellent']}%"
        row.cells[3].text = f"{data['pass']}%"
        for cell in row.cells:
            for paragraph in cell.paragraphs:
                paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
    
    # 总分行
    table.add_row()
    row = table.rows[6]
    row.cells[0].text = '总分'
    row.cells[1].text = str(overall['total']['avg'])
    row.cells[2].text = f"{overall['total']['excellent']}%"
    row.cells[3].text = f"{overall['total']['pass']}%"
    for cell in row.cells:
        for paragraph in cell.paragraphs:
            paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
    
    doc.add_paragraph('')
    
    # 三、各班成绩排名
    doc.add_heading('三、各班成绩排名（按总分平均分）', level=1)
    
    headers = ['排名', '班级', '人数', '≥400分占比', '≥360分占比', '总分平均分', '优秀率', '及格率']
    table = doc.add_table(rows=9, cols=8, style='Light Shading Accent 1')
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    
    for j, header in enumerate(headers):
        cell = table.rows[0].cells[j]
        cell.text = header
        for paragraph in cell.paragraphs:
            paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
    
    for i, cls in enumerate(ranked_classes):
        row = table.rows[i+1]
        row.cells[0].text = str(i+1)
        row.cells[1].text = f"{cls['id']}班"
        row.cells[2].text = str(cls['students'])
        row.cells[3].text = f"{cls['score_400_pct']}%"
        row.cells[4].text = f"{cls['score_360_pct']}%"
        row.cells[5].text = str(cls['total']['avg'])
        row.cells[6].text = f"{cls['total']['excellent']}%"
        row.cells[7].text = f"{cls['total']['pass']}%"
        for cell in row.cells:
            for paragraph in cell.paragraphs:
                paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
    
    doc.add_paragraph('')
    
    # 四、优势与短板分析
    doc.add_heading('四、优势与短板分析', level=1)
    
    # 找出最优和最差科目
    best_subject = max(overall['subjects'].items(), key=lambda x: x[1]['avg'])
    worst_subject = min(overall['subjects'].items(), key=lambda x: x[1]['pass'])
    
    doc.add_heading('优势科目', level=2)
    p = doc.add_paragraph()
    p.add_run(f'{best_subject[0]}').bold = True
    p.add_run(f"：平均分{best_subject[1]['avg']}分，及格率{best_subject[1]['pass']}%")
    
    doc.add_heading('薄弱科目', level=2)
    p = doc.add_paragraph()
    p.add_run(f'{worst_subject[0]}').bold = True
    p.add_run(f"：平均分{worst_subject[1]['avg']}分，及格率仅{worst_subject[1]['pass']}% — ")
    run = p.add_run('⚠️ 需重点关注')
    run.font.color.rgb = RGBColor(0xc0, 0x39, 0x2b)
    
    doc.add_paragraph('')
    
    # 五、班级表现分析
    doc.add_heading('五、班级表现分析', level=1)
    
    doc.add_heading('表现突出的班级', level=2)
    top3 = ranked_classes[:3]
    for cls in top3:
        p = doc.add_paragraph()
        p.add_run(f'{cls["id"]}班').bold = True
        p.add_run(f"：总分平均分{cls['total']['avg']}分，≥400分占比{cls['score_400_pct']}%，≥360分占比{cls['score_360_pct']}%")
    
    doc.add_heading('需要关注的班级', level=2)
    bottom2 = ranked_classes[-2:]
    for cls in bottom2:
        p = doc.add_paragraph()
        p.add_run(f'{cls["id"]}班').bold = True
        p.add_run(f"：总分平均分{cls['total']['avg']}分，≥360分占比仅{cls['score_360_pct']}%")
    
    doc.add_paragraph('')
    
    # 六、后期教学建议
    doc.add_heading('六、后期教学建议', level=1)
    
    suggestions = [
        (f'{worst_subject[0]}学科（最紧急）', [
            f"及格率仅{worst_subject[1]['pass']}%，是最大拉分项",
            "建议分层教学，后20%学生重点抓基础",
            "增加专项训练时间",
            "建立学科互助小组"
        ]),
        ('保持优势学科', [
            f"{best_subject[0]}保持当前教学节奏",
            "总结优秀教学经验，在年级内推广",
            "适当增加拓展内容"
        ]),
        ('尖子生培养', [
            f"400分以上{overall['score_400']}人（{overall['score_400_pct']}%）",
            "组建培优小组，针对性拔高训练",
            "目标：期末400分以上占比提升至15%"
        ]),
        ('后20%学生帮扶', [
            "建立'一生一策'跟踪档案",
            "安排课后辅导，重点抓基础题",
            "家校联动，形成合力"
        ]),
    ]
    
    for title, items in suggestions:
        doc.add_heading(title, level=2)
        for item in items:
            doc.add_paragraph(item, style='List Bullet')
    
    doc.add_paragraph('')
    
    # 七、期末目标
    doc.add_heading('七、期末目标建议', level=1)
    
    table = doc.add_table(rows=5, cols=3, style='Light Shading Accent 1')
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    headers = ['指标', '当前', '期末目标']
    goals = [
        ('总分平均分', f"{overall['total']['avg']}分", '280+'),
        ('英语及格率', f"{overall['subjects']['英语']['pass']}%", '40%+'),
        ('400分以上占比', f"{overall['score_400_pct']}%", '15%'),
        ('360分以上占比', f"{overall['score_360_pct']}%", '35%'),
    ]
    
    for j, header in enumerate(headers):
        table.rows[0].cells[j].text = header
        for paragraph in table.rows[0].cells[j].paragraphs:
            paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
    
    for i, (indicator, current, target) in enumerate(goals):
        row = table.rows[i+1]
        row.cells[0].text = indicator
        row.cells[1].text = current
        row.cells[2].text = target
        for cell in row.cells:
            for paragraph in cell.paragraphs:
                paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
    
    word_path = os.path.join(OUTPUT_DIR, '八年级期中考试成绩分析报告.docx')
    doc.save(word_path)
    print(f"Word文档已生成：{word_path}")
    return word_path


# ==================== PPT 生成 ====================
def create_ppt():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    PRIMARY = PPTColor(0x1a, 0x5c, 0x8a)
    ACCENT = PPTColor(0x2e, 0x86, 0xc1)
    WHITE = PPTColor(0xff, 0xff, 0xff)
    DARK = PPTColor(0x2c, 0x3e, 0x50)
    GREEN = PPTColor(0x27, 0xae, 0x60)
    RED = PPTColor(0xc0, 0x39, 0x2b)
    ORANGE = PPTColor(0xe6, 0x7e, 0x22)
    
    def add_bg(slide, color):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = color
    
    def add_tb(slide, left, top, width, height, text, size=18, color=DARK, bold=False, align=PP_ALIGN.LEFT):
        txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.alignment = align
    
    def add_tbl(slide, left, top, width, height, rows, cols):
        return slide.shapes.add_table(rows, cols, Inches(left), Inches(top), Inches(width), Inches(height)).table
    
    def style_cell(cell, text, size=11, color=DARK, bold=False, align=PP_ALIGN.CENTER, bg=None):
        cell.text = ''
        p = cell.text_frame.paragraphs[0]
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.alignment = align
        if bg:
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
    
    best_subject = max(overall['subjects'].items(), key=lambda x: x[1]['avg'])
    worst_subject = min(overall['subjects'].items(), key=lambda x: x[1]['pass'])
    sub = overall['subjects']
    
    # Slide 1: 封面
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, PRIMARY)
    add_tb(s, 2, 2, 9.333, 1.5, '八年级期中考试成绩分析报告', size=36, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_tb(s, 2, 3.8, 9.333, 0.8, f'参考人数：{overall["total_students"]}人  |  总分平均分：{overall["total"]["avg"]}分', size=18, color=PPTColor(0xbb, 0xdf, 0xf1), align=PP_ALIGN.CENTER)
    add_tb(s, 2, 5, 9.333, 0.6, f'400分以上：{overall["score_400"]}人（{overall["score_400_pct"]}%)', size=20, color=PPTColor(0xf9, 0xe7, 0x9f), bold=True, align=PP_ALIGN.CENTER)
    add_tb(s, 2, 6.2, 9.333, 0.6, '2026年5月', size=14, color=PPTColor(0xbb, 0xdf, 0xf1), align=PP_ALIGN.CENTER)
    
    # Slide 2: 整体概况
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, WHITE)
    add_tb(s, 0.5, 0.3, 12, 0.8, '一、整体概况', size=28, color=PRIMARY, bold=True)
    
    metrics = [
        ('班级数', '8个班', ACCENT),
        ('参考人数', f'{overall["total_students"]}人', PRIMARY),
        ('≥400分', f'{overall["score_400"]}人 ({overall["score_400_pct"]}%)', GREEN),
        ('≥360分', f'{overall["score_360"]}人 ({overall["score_360_pct"]}%)', ORANGE),
    ]
    for i, (label, value, color) in enumerate(metrics):
        left = 0.5 + i * 3.1
        shape = s.shapes.add_shape(1, Inches(left), Inches(1.5), Inches(2.8), Inches(2.2))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        add_tb(s, left+0.2, 1.7, 2.4, 0.6, label, size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_tb(s, left+0.2, 2.3, 2.4, 1.2, value, size=24, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    
    add_tb(s, 0.5, 4.2, 12, 1.5, f'📊 整体表现：\n• 总分平均分{overall["total"]["avg"]}分\n• 优秀率{overall["total"]["excellent"]}%，及格率{overall["total"]["pass"]}%\n• 各班差距较大，最高分与最低分相差较大', size=14, color=DARK)
    
    # Slide 3: 各科成绩详情
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, WHITE)
    add_tb(s, 0.5, 0.3, 12, 0.8, '二、各科成绩详情', size=28, color=PRIMARY, bold=True)
    
    tbl = add_tbl(s, 1, 1.5, 11, 3.5, 7, 5)
    for j, h in enumerate(['科目','平均分','优秀率','及格率','评价']):
        style_cell(tbl.cell(0,j), h, size=14, color=WHITE, bold=True, bg=PRIMARY)
    
    evals = {k: '✅ 优势' if v['avg'] > 60 else ('⚠️ 需关注' if v['pass'] < 40 else '📊 中等') for k, v in sub.items()}
    all_sub = list(sub.items()) + [('总分', overall['total'])]
    for i, (name, d) in enumerate(all_sub):
        r = i+1
        bg = PPTColor(0xd5,0xdb,0xdb) if name=='总分' else None
        style_cell(tbl.cell(r,0), name, size=13, bold=(name=='总分'), bg=bg)
        style_cell(tbl.cell(r,1), str(d['avg']), size=13, bg=bg)
        style_cell(tbl.cell(r,2), f"{d['excellent']}%", size=13, bg=bg)
        if name == worst_subject[0]:
            style_cell(tbl.cell(r,3), f"{d['pass']}%", size=13, color=RED, bold=True, bg=bg)
        else:
            style_cell(tbl.cell(r,3), f"{d['pass']}%", size=13, bg=bg)
        style_cell(tbl.cell(r,4), evals.get(name, '—'), size=12, bg=bg)
    
    add_tb(s, 0.5, 5.5, 12, 1.5, f'📌 关键发现：\n• {best_subject[0]}（{best_subject[1]["avg"]}分）是优势科目\n• {worst_subject[0]}（{worst_subject[1]["avg"]}分）及格率仅{worst_subject[1]["pass"]}%，是最大短板\n• 需重点关注薄弱学科', size=14, color=DARK)
    
    # Slide 4: 班级排名
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, WHITE)
    add_tb(s, 0.5, 0.3, 12, 0.8, '三、各班成绩排名', size=28, color=PRIMARY, bold=True)
    
    tbl = add_tbl(s, 0.8, 1.5, 11.5, 4, 9, 6)
    headers = ['排名','班级','人数','≥400分','≥360分','总分平均分']
    for j, h in enumerate(headers):
        style_cell(tbl.cell(0,j), h, size=13, color=WHITE, bold=True, bg=PRIMARY)
    
    for i, cls in enumerate(ranked_classes):
        r = i+1
        is_top3 = i < 3
        is_bottom2 = i >= 6
        bg = PPTColor(0xd4,0xef,0xdf) if is_top3 else (PPTColor(0xff,0xe6,0xe6) if is_bottom2 else None)
        tc = PRIMARY if is_top3 else (RED if is_bottom2 else DARK)
        b = is_top3 or is_bottom2
        style_cell(tbl.cell(r,0), str(i+1), size=12, color=tc, bold=b, bg=bg)
        style_cell(tbl.cell(r,1), f"{cls['id']}班", size=12, color=tc, bold=b, bg=bg)
        style_cell(tbl.cell(r,2), str(cls['students']), size=12, bg=bg)
        style_cell(tbl.cell(r,3), f"{cls['score_400_pct']}%", size=12, bg=bg)
        style_cell(tbl.cell(r,4), f"{cls['score_360_pct']}%", size=12, bg=bg)
        style_cell(tbl.cell(r,5), str(cls['total']['avg']), size=12, color=tc, bold=b, bg=bg)
    
    add_tb(s, 0.5, 5.8, 12, 1.2, f'📊 排名分析：\n• 第1名：{ranked_classes[0]["id"]}班（{ranked_classes[0]["total"]["avg"]}分）\n• 最后1名：{ranked_classes[-1]["id"]}班（{ranked_classes[-1]["total"]["avg"]}分）\n• 班级差距明显，需关注落后班级', size=13, color=DARK)
    
    # Slide 5: 优势与短板
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, WHITE)
    add_tb(s, 0.5, 0.3, 12, 0.8, '四、优势与短板分析', size=28, color=PRIMARY, bold=True)
    
    add_tb(s, 0.5, 1.3, 6, 0.6, '✅ 优势科目', size=22, color=GREEN, bold=True)
    add_tb(s, 0.5, 2.0, 6, 0.5, f'{best_subject[0]}：平均分{best_subject[1]["avg"]}分  |  及格率{best_subject[1]["pass"]}%', size=16, color=DARK, bold=True)
    add_tb(s, 0.5, 2.6, 6, 1.5, '• 平均分最高\n• 教学基础扎实\n• 保持当前教学节奏', size=13, color=DARK)
    
    add_tb(s, 7, 1.3, 6, 0.6, '⚠️ 薄弱科目', size=22, color=RED, bold=True)
    add_tb(s, 7, 2.0, 6, 0.5, f'{worst_subject[0]}：平均分{worst_subject[1]["avg"]}分  |  及格率{worst_subject[1]["pass"]}%', size=16, color=RED, bold=True)
    add_tb(s, 7, 2.6, 6, 1.5, '• 最大短板，需重点关注\n• 及格率远低于其他科目\n• 建议分层教学，抓基础', size=13, color=DARK)
    
    # Slide 6: 教学建议
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, WHITE)
    add_tb(s, 0.5, 0.3, 12, 0.8, '五、后期教学建议', size=28, color=PRIMARY, bold=True)
    
    suggestions = [
        (f'1️⃣ {worst_subject[0]}学科（最紧急）', [f'及格率仅{worst_subject[1]["pass"]}%，是最大拉分项','分层教学，后20%学生重点抓基础','增加专项训练时间','建立学科互助小组'], RED),
        ('2️⃣ 保持优势', [f'{best_subject[0]}保持当前教学节奏','总结优秀教学经验，在年级内推广','适当增加拓展内容'], GREEN),
        ('3️⃣ 尖子生培养', [f'400分以上{overall["score_400"]}人（{overall["score_400_pct"]}%)','组建培优小组，针对性拔高训练','目标：期末400分以上占比提升至15%'], ACCENT),
        ('4️⃣ 后20%帮扶', ['建立"一生一策"跟踪档案','安排课后辅导，重点抓基础题','家校联动，设置进步奖'], PRIMARY),
    ]
    
    y = 1.3
    for title, items, color in suggestions:
        add_tb(s, 0.5, y, 12, 0.5, title, size=16, color=color, bold=True)
        y += 0.5
        for item in items:
            add_tb(s, 0.8, y, 11.5, 0.35, f'• {item}', size=12, color=DARK)
            y += 0.35
        y += 0.15
    
    # Slide 7: 期末目标
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, WHITE)
    add_tb(s, 0.5, 0.3, 12, 0.8, '六、期末目标建议', size=28, color=PRIMARY, bold=True)
    
    tbl = add_tbl(s, 2, 1.5, 9, 3, 5, 4)
    for j, h in enumerate(['指标','当前值','期末目标','提升幅度']):
        style_cell(tbl.cell(0,j), h, size=14, color=WHITE, bold=True, bg=PRIMARY)
    
    goals = [
        ('总分平均分', f'{overall["total"]["avg"]}分', '280+', '+5+'),
        (f'{worst_subject[0]}及格率', f'{worst_subject[1]["pass"]}%', '40%+', '+10+'),
        ('400分以上占比', f'{overall["score_400_pct"]}%', '15%', '+5+'),
        ('360分以上占比', f'{overall["score_360_pct"]}%', '35%', '+5+'),
    ]
    for i, (ind, cur, tgt, imp) in enumerate(goals):
        r = i+1
        style_cell(tbl.cell(r,0), ind, size=14, bold=True)
        style_cell(tbl.cell(r,1), cur, size=14)
        style_cell(tbl.cell(r,2), tgt, size=14, color=GREEN, bold=True)
        style_cell(tbl.cell(r,3), imp, size=14, color=ACCENT)
    
    add_tb(s, 0.5, 5.2, 12, 2, '🎯 核心目标：\n• 总分平均分突破280分\n• 薄弱学科及格率提升至40%以上\n• 400分以上占比提升至15%\n• 缩小班级差距，整体提升', size=16, color=DARK, bold=True, align=PP_ALIGN.CENTER)
    
    ppt_path = os.path.join(OUTPUT_DIR, '八年级期中考试成绩分析汇报.pptx')
    prs.save(ppt_path)
    print(f"PPT已生成：{ppt_path}")
    return ppt_path


if __name__ == '__main__':
    word_path = create_word_report()
    ppt_path = create_ppt()
    print(f'\n✅ 文件生成完成！')
    print(f'Word: {word_path}')
    print(f'PPT: {ppt_path}')
