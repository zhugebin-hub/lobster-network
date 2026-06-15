#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
科技感计算机专业教学 PPT 模板生成器
"""

from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# 创建演示文稿
prs = Presentation()
prs.slide_width = Cm(33.867)
prs.slide_height = Cm(19.05)

# 科技感颜色
COLORS = {
    'bg_dark': (10, 15, 30),          # 深空蓝黑
    'bg_grad1': (20, 30, 60),         # 渐变蓝
    'bg_grad2': (5, 10, 25),          # 渐变深蓝
    'neon_blue': (0, 210, 255),       # 霓虹蓝
    'neon_cyan': (0, 255, 255),       # 青色
    'neon_purple': (138, 43, 226),    # 紫色
    'neon_green': (0, 255, 127),      # 绿色
    'text_bright': (255, 255, 255),   # 亮白
    'text_dim': (150, 170, 200),      # 暗淡文字
    'accent_gold': (255, 215, 0),     # 金色
}

def set_background(slide, color=COLORS['bg_dark']):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*color)

def add_tech_grid(slide):
    """添加科技网格背景"""
    # 水平线
    for i in range(12):
        y = Cm(i * 1.6)
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Cm(0), y, prs.slide_width, Cm(0.02)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*COLORS['neon_blue'])
        shape.fill.transparency = 0.85
        shape.line.fill.background()
    
    # 垂直线
    for i in range(22):
        x = Cm(i * 1.55)
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, x, Cm(0), Cm(0.02), prs.slide_height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*COLORS['neon_blue'])
        shape.fill.transparency = 0.85
        shape.line.fill.background()

def add_circuit_lines(slide, start_x, start_y, direction='right'):
    """添加电路线条装饰"""
    # 简化版电路线
    for i in range(3):
        x = start_x + Cm(i * 2)
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, x, start_y, Cm(1.5), Cm(0.05)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*COLORS['neon_cyan'])
        shape.fill.transparency = 0.6
        shape.line.fill.background()

def add_hexagon(slide, x, y, size=Cm(2), color=COLORS['neon_blue'], transparency=0.7):
    """添加六边形装饰"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.HEXAGON, x, y, size, size * 0.9
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*color)
    shape.fill.transparency = transparency
    shape.line.color.rgb = RGBColor(*color)
    shape.line.width = Pt(1)

def add_header(slide, title, subtitle=""):
    # 顶部科技条
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Cm(0), Cm(0), prs.slide_width, Cm(0.4)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*COLORS['neon_blue'])
    shape.line.fill.background()
    
    # 装饰光点
    for i in range(8):
        x = Cm(2 + i * 4)
        shape = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, x, Cm(0.1), Cm(0.3), Cm(0.3)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*COLORS['neon_cyan'])
        shape.fill.transparency = 0.5
        shape.line.fill.background()
    
    # 标题
    textbox = slide.shapes.add_textbox(Cm(1.5), Cm(0.6), Cm(25), Cm(1))
    tf = textbox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = RGBColor(*COLORS['text_bright'])
    p.font.name = 'Microsoft YaHei'
    
    if subtitle:
        textbox = slide.shapes.add_textbox(Cm(22), Cm(0.7), Cm(10), Cm(0.8))
        tf = textbox.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(*COLORS['neon_cyan'])
        p.font.name = 'Microsoft YaHei'
        p.alignment = PP_ALIGN.RIGHT

def add_footer(slide, page_num=""):
    # 底线
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Cm(2), prs.slide_height - Cm(1), prs.slide_width - Cm(4), Cm(0.05)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*COLORS['neon_blue'])
    shape.fill.transparency = 0.5
    shape.line.fill.background()
    
    # 页码
    textbox = slide.shapes.add_textbox(prs.slide_width - Cm(3), prs.slide_height - Cm(1.3), Cm(2), Cm(0.5))
    tf = textbox.text_frame
    p = tf.paragraphs[0]
    p.text = page_num
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(*COLORS['text_dim'])
    p.font.name = 'Consolas'
    p.alignment = PP_ALIGN.RIGHT
    
    # 装饰六边形
    add_hexagon(slide, Cm(2.5), prs.slide_height - Cm(1.5), Cm(0.6), COLORS['neon_purple'])

def add_title_slide(prs, title, subtitle, footer=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, COLORS['bg_dark'])
    
    # 背景网格
    add_tech_grid(slide)
    
    # 中央光晕效果
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Cm(10), Cm(5), Cm(14), Cm(9)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*COLORS['neon_blue'])
    shape.fill.transparency = 0.9
    shape.line.fill.background()
    
    # 标题框
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Cm(5), Cm(6), prs.slide_width - Cm(10), Cm(5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*COLORS['bg_grad1'])
    shape.fill.transparency = 0.5
    shape.line.color.rgb = RGBColor(*COLORS['neon_blue'])
    shape.line.width = Pt(2)
    
    # 主标题 - 渐变效果模拟
    textbox = slide.shapes.add_textbox(Cm(6), Cm(7), prs.slide_width - Cm(12), Cm(2.5))
    tf = textbox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(52)
    p.font.bold = True
    p.font.color.rgb = RGBColor(*COLORS['text_bright'])
    p.font.name = 'Microsoft YaHei'
    p.alignment = PP_ALIGN.CENTER
    
    # 副标题
    textbox = slide.shapes.add_textbox(Cm(6), Cm(9.5), prs.slide_width - Cm(12), Cm(1.5))
    tf = textbox.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(*COLORS['neon_cyan'])
    p.font.name = 'Microsoft YaHei'
    p.alignment = PP_ALIGN.CENTER
    
    # 底部信息
    if footer:
        textbox = slide.shapes.add_textbox(Cm(6), prs.slide_height - Cm(3), prs.slide_width - Cm(12), Cm(1))
        tf = textbox.text_frame
        p = tf.paragraphs[0]
        p.text = footer
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(*COLORS['text_dim'])
        p.font.name = 'Consolas'
        p.alignment = PP_ALIGN.CENTER
    
    # 角落装饰
    add_hexagon(slide, Cm(1), Cm(1), Cm(1.5), COLORS['neon_purple'], 0.6)
    add_hexagon(slide, prs.slide_width - Cm(2.5), Cm(1), Cm(1.5), COLORS['neon_green'], 0.6)
    add_hexagon(slide, Cm(1), prs.slide_height - Cm(2.5), Cm(1.5), COLORS['neon_cyan'], 0.6)
    add_hexagon(slide, prs.slide_width - Cm(2.5), prs.slide_height - Cm(2.5), Cm(1.5), COLORS['neon_blue'], 0.6)
    
    return slide

def add_code_slide(prs, title, code_text, page_num=""):
    """代码展示页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, COLORS['bg_dark'])
    add_header(slide, title)
    
    # 代码框
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Cm(2), Cm(2.5), prs.slide_width - Cm(4), Cm(12)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*COLORS['bg_grad2'])
    shape.line.color.rgb = RGBColor(*COLORS['neon_blue'])
    shape.line.width = Pt(2)
    
    # 代码文字
    textbox = slide.shapes.add_textbox(Cm(3), Cm(3.5), prs.slide_width - Cm(6), Cm(10))
    tf = textbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = code_text
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(*COLORS['neon_green'])
    p.font.name = 'Consolas'
    
    add_footer(slide, page_num)
    return slide

def add_content_slide(prs, title, content_items, page_num=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, COLORS['bg_dark'])
    add_header(slide, title)
    
    # 背景装饰
    add_tech_grid(slide)
    
    y = Cm(3)
    for item in content_items:
        if isinstance(item, dict):
            if item.get('type') == 'title':
                textbox = slide.shapes.add_textbox(Cm(2), y, Cm(30), Cm(0.8))
                tf = textbox.text_frame
                p = tf.paragraphs[0]
                p.text = "► " + item['text']
                p.font.size = Pt(22)
                p.font.bold = True
                p.font.color.rgb = RGBColor(*COLORS['neon_cyan'])
                p.font.name = 'Microsoft YaHei'
                y += Cm(1.2)
            elif item.get('type') == 'bullet':
                textbox = slide.shapes.add_textbox(Cm(3), y, Cm(28), Cm(0.7))
                tf = textbox.text_frame
                p = tf.paragraphs[0]
                p.text = "• " + item['text']
                p.font.size = Pt(18)
                p.font.color.rgb = RGBColor(*COLORS['text_bright'])
                p.font.name = 'Microsoft YaHei'
                y += Cm(0.9)
            elif item.get('type') == 'highlight':
                shape = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE, Cm(2), y, Cm(30), Cm(1.5)
                )
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor(*COLORS['neon_blue'])
                shape.fill.transparency = 0.8
                shape.line.color.rgb = RGBColor(*COLORS['neon_blue'])
                shape.line.width = Pt(2)
                
                textbox = slide.shapes.add_textbox(Cm(3), y + Cm(0.3), Cm(28), Cm(0.9))
                tf = textbox.text_frame
                p = tf.paragraphs[0]
                p.text = item['text']
                p.font.size = Pt(18)
                p.font.bold = True
                p.font.color.rgb = RGBColor(*COLORS['neon_cyan'])
                p.font.name = 'Microsoft YaHei'
                y += Cm(2)
            elif item.get('type') == 'box':
                shape = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE, Cm(2), y, Cm(30), item.get('height', Cm(2))
                )
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor(*COLORS['bg_grad2'])
                shape.line.color.rgb = RGBColor(*COLORS['neon_purple'])
                shape.line.width = Pt(2)
                
                textbox = slide.shapes.add_textbox(Cm(3), y + Cm(0.3), Cm(28), Cm(1.5))
                tf = textbox.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = item['text']
                p.font.size = Pt(16)
                p.font.color.rgb = RGBColor(*COLORS['text_dim'])
                p.font.name = 'Microsoft YaHei'
                y += item.get('height', Cm(2)) + Cm(0.3)
            elif item.get('type') == 'tech_card':
                # 科技卡片
                card_w = item.get('width', Cm(13))
                card_h = item.get('height', Cm(4))
                card_x = item.get('x', Cm(2))
                
                shape = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE, card_x, y, card_w, card_h
                )
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor(*COLORS['bg_grad1'])
                shape.fill.transparency = 0.7
                shape.line.color.rgb = RGBColor(*COLORS.get(item.get('color', 'neon_blue'), COLORS['neon_blue']))
                shape.line.width = Pt(2)
                
                # 卡片标题
                textbox = slide.shapes.add_textbox(card_x + Cm(0.5), y + Cm(0.3), card_w - Cm(1), Cm(0.6))
                tf = textbox.text_frame
                p = tf.paragraphs[0]
                p.text = item.get('card_title', '')
                p.font.size = Pt(16)
                p.font.bold = True
                p.font.color.rgb = RGBColor(*COLORS['neon_cyan'])
                p.font.name = 'Microsoft YaHei'
                
                # 卡片内容
                textbox = slide.shapes.add_textbox(card_x + Cm(0.5), y + Cm(1), card_w - Cm(1), card_h - Cm(1.5))
                tf = textbox.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = item.get('card_content', '')
                p.font.size = Pt(14)
                p.font.color.rgb = RGBColor(*COLORS['text_bright'])
                p.font.name = 'Microsoft YaHei'
                
                y += card_h + Cm(0.3)
        else:
            textbox = slide.shapes.add_textbox(Cm(3), y, Cm(28), Cm(0.7))
            tf = textbox.text_frame
            p = tf.paragraphs[0]
            p.text = item
            p.font.size = Pt(18)
            p.font.color.rgb = RGBColor(*COLORS['text_bright'])
            p.font.name = 'Microsoft YaHei'
            y += Cm(0.9)
    
    add_footer(slide, page_num)
    return slide

def add_diagram_slide(prs, title, diagram_type, page_num=""):
    """架构图页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, COLORS['bg_dark'])
    add_header(slide, title)
    
    # 背景网格
    add_tech_grid(slide)
    
    if diagram_type == 'layers':
        # 分层架构图
        layers = [
            (Cm(5), Cm(4), Cm(24), Cm(1.8), "应用层", "Web/App/小程序"),
            (Cm(5), Cm(6.5), Cm(24), Cm(1.8), "服务层", "API 网关/微服务"),
            (Cm(5), Cm(9), Cm(24), Cm(1.8), "数据层", "MySQL/MongoDB/Redis"),
            (Cm(5), Cm(11.5), Cm(24), Cm(1.8), "基础设施", "云服务器/容器/CDN"),
        ]
        
        colors = [COLORS['neon_purple'], COLORS['neon_blue'], COLORS['neon_cyan'], COLORS['neon_green']]
        
        for i, (x, y, w, h, title_text, content) in enumerate(layers):
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, x, y, w, h
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(*colors[i])
            shape.fill.transparency = 0.7
            shape.line.color.rgb = RGBColor(*colors[i])
            shape.line.width = Pt(2)
            
            textbox = slide.shapes.add_textbox(x + Cm(1), y + Cm(0.3), w - Cm(2), Cm(0.6))
            tf = textbox.text_frame
            p = tf.paragraphs[0]
            p.text = title_text
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = RGBColor(*COLORS['text_bright'])
            p.font.name = 'Microsoft YaHei'
            
            textbox = slide.shapes.add_textbox(x + Cm(1), y + Cm(1), w - Cm(2), Cm(0.6))
            tf = textbox.text_frame
            p = tf.paragraphs[0]
            p.text = content
            p.font.size = Pt(14)
            p.font.color.rgb = RGBColor(*COLORS['text_dim'])
            p.font.name = 'Microsoft YaHei'
    
    add_footer(slide, page_num)
    return slide

def add_comparison_slide(prs, title, headers, rows, page_num=""):
    """对比表格页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, COLORS['bg_dark'])
    add_header(slide, title)
    
    # 背景网格
    add_tech_grid(slide)
    
    col_width = Cm(6.5)
    row_height = Cm(1.5)
    start_x = Cm(2)
    start_y = Cm(3)
    
    # 表头
    for i, header in enumerate(headers):
        x = start_x + i * col_width
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, x, start_y, col_width, Cm(1.2)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*COLORS['neon_blue'])
        shape.fill.transparency = 0.7
        shape.line.color.rgb = RGBColor(*COLORS['neon_blue'])
        shape.line.width = Pt(2)
        
        textbox = slide.shapes.add_textbox(x + Cm(0.3), start_y + Cm(0.2), col_width - Cm(0.6), Cm(0.8))
        tf = textbox.text_frame
        p = tf.paragraphs[0]
        p.text = header
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = RGBColor(*COLORS['text_bright'])
        p.font.name = 'Microsoft YaHei'
        p.alignment = PP_ALIGN.CENTER
    
    # 内容行
    for row_idx, row_data in enumerate(rows):
        y = start_y + Cm(1.2) + row_idx * row_height
        for col_idx, cell in enumerate(row_data):
            x = start_x + col_idx * col_width
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, x, y, col_width - Cm(0.1), row_height - Cm(0.1)
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(*COLORS['bg_grad2'])
            shape.line.color.rgb = RGBColor(*COLORS['neon_blue'])
            shape.line.width = Pt(1)
            
            textbox = slide.shapes.add_textbox(x + Cm(0.3), y + Cm(0.3), col_width - Cm(0.9), row_height - Cm(0.6))
            tf = textbox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = cell
            p.font.size = Pt(12)
            p.font.color.rgb = RGBColor(*COLORS['text_bright'])
            p.font.name = 'Microsoft YaHei'
    
    add_footer(slide, page_num)
    return slide

def add_end_slide(prs, title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, COLORS['bg_dark'])
    
    # 背景网格
    add_tech_grid(slide)
    
    # 中央光晕
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Cm(10), Cm(5), Cm(14), Cm(9)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*COLORS['neon_blue'])
    shape.fill.transparency = 0.9
    shape.line.fill.background()
    
    # 标题
    textbox = slide.shapes.add_textbox(Cm(6), Cm(7), prs.slide_width - Cm(12), Cm(3))
    tf = textbox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(50)
    p.font.bold = True
    p.font.color.rgb = RGBColor(*COLORS['neon_cyan'])
    p.font.name = 'Microsoft YaHei'
    p.alignment = PP_ALIGN.CENTER
    
    if subtitle:
        textbox = slide.shapes.add_textbox(Cm(6), Cm(10), prs.slide_width - Cm(12), Cm(2))
        tf = textbox.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(22)
        p.font.color.rgb = RGBColor(*COLORS['text_dim'])
        p.font.name = 'Microsoft YaHei'
        p.alignment = PP_ALIGN.CENTER
    
    # 角落装饰
    add_hexagon(slide, Cm(1), Cm(1), Cm(1.5), COLORS['neon_purple'], 0.6)
    add_hexagon(slide, prs.slide_width - Cm(2.5), prs.slide_height - Cm(2.5), Cm(1.5), COLORS['neon_green'], 0.6)
    
    return slide

# ============ 生成模板幻灯片 ============

# 1. 封面
add_title_slide(prs,
    title="计算机专业技术",
    subtitle="科技感教学 PPT 模板",
    footer="适用于编程/算法/系统架构等课程"
)

# 2. 目录
add_content_slide(prs, "课程目录", [
    {"type": "title", "text": "CONTENT"},
    "",
    {"type": "bullet", "text": "01 技术概述与技术栈介绍"},
    {"type": "bullet", "text": "02 核心原理与算法分析"},
    {"type": "bullet", "text": "03 系统架构与设计模式"},
    {"type": "bullet", "text": "04 实战案例与代码演示"},
    {"type": "bullet", "text": "05 性能优化与最佳实践"},
    {"type": "bullet", "text": "06 课后实训与拓展任务"},
], "01/12")

# 3. 技术概述
add_content_slide(prs, "技术概述", [
    {"type": "highlight", "text": "技术定义与核心概念"},
    "",
    {"type": "box", "text": "在此处输入技术的定义和核心概念说明...", "height": Cm(2)},
    "",
    {"type": "title", "text": "技术特点"},
    {"type": "bullet", "text": "高性能 - 处理效率高，响应速度快"},
    {"type": "bullet", "text": "可扩展 - 支持水平扩展，应对高并发"},
    {"type": "bullet", "text": "易维护 - 代码结构清晰，文档完善"},
    {"type": "bullet", "text": "安全可靠 - 多重防护，数据加密"},
], "02/12")

# 4. 技术对比
add_comparison_slide(prs, "技术方案对比",
    ["方案", "性能", "复杂度", "适用场景", "学习曲线"],
    [
        ["方案 A", "★★★★★", "中等", "高并发场景", "平缓"],
        ["方案 B", "★★★★☆", "较低", "快速开发", "简单"],
        ["方案 C", "★★★☆☆", "较高", "定制化需求", "陡峭"],
        ["方案 D", "★★★★☆", "中等", "企业级应用", "中等"],
    ],
    "03/12"
)

# 5. 系统架构
add_diagram_slide(prs, "系统架构图", "layers", "04/12")

# 6. 核心原理
add_content_slide(prs, "核心原理", [
    {"type": "highlight", "text": "关键技术原理详解"},
    "",
    {"type": "title", "text": "原理一"},
    {"type": "box", "text": "在此处详细描述第一个核心原理...", "height": Cm(2)},
    "",
    {"type": "title", "text": "原理二"},
    {"type": "box", "text": "在此处详细描述第二个核心原理...", "height": Cm(2)},
    "",
    {"type": "title", "text": "原理三"},
    {"type": "box", "text": "在此处详细描述第三个核心原理...", "height": Cm(2)},
], "05/12")

# 7. 代码示例
code_example = """// 示例代码：核心算法实现

def core_algorithm(data):
    \"\"\"
    核心算法函数
    参数：data - 输入数据
    返回：处理结果
    \"\"\"
    # 数据预处理
    processed = preprocess(data)
    
    # 核心逻辑
    result = []
    for item in processed:
        if validate(item):
            result.append(transform(item))
    
    # 返回结果
    return optimize(result)

# 时间复杂度：O(n log n)
# 空间复杂度：O(n)"""

add_code_slide(prs, "代码示例", code_example, "06/12")

# 8. 算法分析
add_content_slide(prs, "算法分析", [
    {"type": "highlight", "text": "时间复杂度与空间复杂度分析"},
    "",
    {"type": "title", "text": "最好情况"},
    {"type": "bullet", "text": "时间复杂度：O(n)"},
    {"type": "bullet", "text": "空间复杂度：O(1)"},
    {"type": "bullet", "text": "适用场景：数据已基本有序"},
    "",
    {"type": "title", "text": "最坏情况"},
    {"type": "bullet", "text": "时间复杂度：O(n²)"},
    {"type": "bullet", "text": "空间复杂度：O(n)"},
    {"type": "bullet", "text": "适用场景：数据完全逆序"},
    "",
    {"type": "title", "text": "平均情况"},
    {"type": "bullet", "text": "时间复杂度：O(n log n)"},
    {"type": "bullet", "text": "空间复杂度：O(log n)"},
], "07/12")

# 9. 实战案例
add_content_slide(prs, "实战案例", [
    {"type": "tech_card", "x": Cm(2), "width": Cm(14), "height": Cm(4), 
     "card_title": "案例一：电商平台", 
     "card_content": "• 场景：秒杀系统\n• 挑战：高并发、防超卖\n• 方案：Redis 缓存 + 队列\n• 成果：QPS 提升 10 倍", 
     "color": "neon_blue"},
    
    {"type": "tech_card", "x": Cm(17), "width": Cm(14), "height": Cm(4), 
     "card_title": "案例二：社交应用", 
     "card_content": "• 场景：消息推送\n• 挑战：实时性、可靠性\n• 方案：WebSocket + MQTT\n• 成果：延迟<100ms", 
     "color": "neon_purple"},
], "08/12")

# 10. 性能优化
add_content_slide(prs, "性能优化", [
    {"type": "highlight", "text": "性能优化策略与方法"},
    "",
    {"type": "title", "text": "数据库优化"},
    {"type": "bullet", "text": "索引优化 - 建立合适索引，避免全表扫描"},
    {"type": "bullet", "text": "查询优化 - 减少 JOIN，使用覆盖索引"},
    {"type": "bullet", "text": "缓存策略 - Redis 缓存热点数据"},
    "",
    {"type": "title", "text": "代码优化"},
    {"type": "bullet", "text": "算法优化 - 选择合适的数据结构和算法"},
    {"type": "bullet", "text": "异步处理 - 耗时操作异步化"},
    {"type": "bullet", "text": "并发控制 - 合理使用线程池"},
    "",
    {"type": "title", "text": "架构优化"},
    {"type": "bullet", "text": "负载均衡 - 分散请求压力"},
    {"type": "bullet", "text": "CDN 加速 - 静态资源就近访问"},
], "09/12")

# 11. 最佳实践
add_content_slide(prs, "最佳实践", [
    {"type": "highlight", "text": "开发规范与最佳实践"},
    "",
    {"type": "title", "text": "代码规范"},
    {"type": "bullet", "text": "遵循命名规范，变量/函数见名知意"},
    {"type": "bullet", "text": "代码注释，关键逻辑添加说明"},
    {"type": "bullet", "text": "单元测试，核心功能覆盖测试"},
    "",
    {"type": "title", "text": "安全实践"},
    {"type": "bullet", "text": "输入验证，防止 SQL 注入/XSS"},
    {"type": "bullet", "text": "权限控制，最小权限原则"},
    {"type": "bullet", "text": "数据加密，敏感信息加密存储"},
    "",
    {"type": "title", "text": "运维实践"},
    {"type": "bullet", "text": "日志记录，便于问题排查"},
    {"type": "bullet", "text": "监控告警，及时发现异常"},
], "10/12")

# 12. 课后实训
add_content_slide(prs, "课后实训", [
    {"type": "highlight", "text": "实训任务与拓展学习"},
    "",
    {"type": "title", "text": "基础任务"},
    {"type": "bullet", "text": "完成课堂代码的复现与调试"},
    {"type": "bullet", "text": "编写技术文档，总结核心知识点"},
    {"type": "bullet", "text": "完成课后练习题"},
    "",
    {"type": "title", "text": "提高任务"},
    {"type": "bullet", "text": "基于所学技术完成一个小项目"},
    {"type": "bullet", "text": "性能对比实验，分析优化效果"},
    {"type": "bullet", "text": "阅读相关技术论文/源码"},
    "",
    {"type": "title", "text": "拓展资源"},
    {"type": "bullet", "text": "官方文档、GitHub 开源项目、技术博客"},
], "11/12")

# 13. 结束页
add_end_slide(prs, "感谢观看", "Q & A 欢迎提问交流")

# 保存文件
output_path = '/home/admin/.openclaw/workspace/科技感计算机专业 PPT 模板.pptx'
prs.save(output_path)
print(f"PPT 模板已生成：{output_path}")
print(f"共 {len(prs.slides)} 页")
