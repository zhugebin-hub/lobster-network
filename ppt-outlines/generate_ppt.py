#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Generate PPT for Christian Sinicization presentation"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from lxml import etree

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color scheme - dark blue + gold, NO red
DARK_BLUE = RGBColor(0x1A, 0x3C, 0x6E)
LIGHT_BLUE = RGBColor(0x2E, 0x86, 0xC1)
GOLD = RGBColor(0xD4, 0xA5, 0x37)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x33, 0x33, 0x33)
GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY = RGBColor(0xF0, 0xF4, 0xF8)

TOTAL = 10
NS_A = 'http://schemas.openxmlformats.org/drawingml/2006/main'

def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, w, h, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text(slide, left, top, w, h, text, size=18, color=BLACK, bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(left, top, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = 'Microsoft YaHei'
    p.alignment = align
    p.space_after = Pt(0)
    p.line_spacing = Pt(int(size * 1.5))
    return tb

def add_bullets(slide, left, top, w, h, items, size=18, color=BLACK):
    tb = slide.shapes.add_textbox(left, top, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = 'Microsoft YaHei'
        p.space_after = Pt(8)
        p.line_spacing = Pt(int(size * 1.6))
        pPr = p._p.find('{' + NS_A + '}pPr')
        if pPr is None:
            pPr = etree.SubElement(p._p, '{%s}pPr' % NS_A)
        buChar = etree.SubElement(pPr, '{%s}buChar' % NS_A)
        buChar.set('char', '\u25b8')
        buClr = etree.SubElement(pPr, '{%s}buClr' % NS_A)
        srgb = etree.SubElement(buClr, '{%s}srgbClr' % NS_A)
        srgb.set('val', '2E86C1')
    return tb

def add_page_num(slide, num):
    add_text(slide, Inches(12.2), Inches(7.0), Inches(1), Inches(0.4),
             f"{num}/{TOTAL}", size=10, color=GRAY, align=PP_ALIGN.RIGHT)

# ==================== Slide 1: Cover ====================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, DARK_BLUE)
add_rect(s, Inches(0), Inches(0), Inches(13.333), Inches(0.08), GOLD)
add_rect(s, Inches(1), Inches(2.8), Inches(3), Inches(0.04), GOLD)

add_text(s, Inches(1), Inches(3.0), Inches(11.333), Inches(1.5),
         '\u4f9d\u6cd5\u800c\u6cbb  \u5411\u5fc3\u800c\u884c', size=44, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(s, Inches(1), Inches(4.5), Inches(11.333), Inches(1),
         '\u2014\u2014\u575a\u6301\u57fa\u7763\u6559\u4e2d\u56fd\u5316\u65b9\u5411\u2014\u2014', size=24, color=GOLD, align=PP_ALIGN.CENTER)
add_text(s, Inches(1), Inches(5.8), Inches(11.333), Inches(0.8),
         '\u5ba3\u8bb2\u4eba\uff1a\u4e00\u7f15\u9633\u5149    2026\u5e745\u6708', size=16, color=WHITE, align=PP_ALIGN.CENTER)
add_page_num(s, 1)

# ==================== Slide 2: Theme ====================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, WHITE)
add_rect(s, Inches(0), Inches(0), Inches(13.333), Inches(0.08), DARK_BLUE)
add_rect(s, Inches(0), Inches(7.3), Inches(13.333), Inches(0.2), LIGHT_BLUE)

add_text(s, Inches(0.8), Inches(0.5), Inches(11.733), Inches(0.8),
         '\u5ba3\u8bb2\u4e3b\u9898', size=32, color=DARK_BLUE, bold=True)
add_rect(s, Inches(0.8), Inches(1.2), Inches(2), Inches(0.06), GOLD)
add_text(s, Inches(0.8), Inches(1.6), Inches(11.733), Inches(0.6),
         '\u4fe1\u4ef0\u81ea\u7531\u4e0e\u9075\u7eaa\u5b88\u6cd5\uff0c\u5982\u4f55\u5728\u5b9e\u8df5\u4e2d\u7edf\u4e00\uff1f', size=22, color=DARK_BLUE, bold=True)

items = [
    '\u56fd\u5bb6\u5b97\u6559\u4e8b\u52a1\u6cd5\u89c4\uff1a\u793e\u4f1a\u79e9\u5e8f\u7684\u4fdd\u969c\uff0c\u66f4\u662f\u4fe1\u4ef0\u5e73\u5b89\u7684\u6839\u57fa',
    '\u201c\u5728\u4e0a\u6709\u6743\u957f\u7684\uff0c\u4eba\u4eba\u5f53\u987a\u670d\u4ed6\u201d\uff08\u7f5713:1\uff09',
    '\u4f9d\u6cd5\u5f00\u5c55\u6d3b\u52a8\uff0c\u662f\u4fe1\u4ef0\u5728\u5f53\u4ee3\u4e2d\u56fd\u6700\u5065\u5eb7\u3001\u6700\u5408\u5b9c\u7684\u8868\u8fbe',
    '\u4e94\u4e2a\u7ef4\u5ea6\uff1a\u4eba\u3001\u5730\u3001\u9053\u3001\u4e8b\u3001\u6cd5'
]
add_bullets(s, Inches(0.8), Inches(2.5), Inches(11.733), Inches(4.5), items, size=18)
add_page_num(s, 2)

# ==================== Slide 3: People ====================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, WHITE)
add_rect(s, Inches(0), Inches(0), Inches(13.333), Inches(0.08), DARK_BLUE)
add_rect(s, Inches(0), Inches(7.3), Inches(13.333), Inches(0.2), LIGHT_BLUE)

add_text(s, Inches(0.8), Inches(0.5), Inches(11.733), Inches(0.8),
         '\u4e00\u3001\u7d27\u62f6\u201c\u4eba\u201d\u7684\u8d44\u8d28', size=32, color=DARK_BLUE, bold=True)
add_rect(s, Inches(0.8), Inches(1.2), Inches(2), Inches(0.06), GOLD)
add_text(s, Inches(0.8), Inches(1.6), Inches(11.733), Inches(0.6),
         '\u843d\u5b9e\u5907\u6848\u5236\uff0c\u786e\u4fdd\u8bb2\u53f0\u6709\u5e8f', size=22, color=DARK_BLUE, bold=True)

items = [
    '\u5b97\u6559\u6d3b\u52a8\u987b\u7531\u4f9d\u6cd5\u8ba4\u5b9a\u5907\u6848\u3001\u6301\u8bc1\u4e0a\u5c97\u7684\u6559\u804c\u4eba\u5458\u4e3b\u6301',
    '\u8de8\u533a\u57df\u8bb2\u9053\u987b\u7ecf\u53cc\u65b9\u56e2\u4f53\u540c\u610f\u5e76\u63d0\u524d\u62a5\u5907',
    '\u4e25\u7981\u672a\u7ecf\u5ba1\u6838\u3001\u8eab\u4efd\u4e0d\u6e05\u3001\u795e\u5b66\u4e0d\u660e\u7684\u4eba\u5458\u767b\u53f0',
    '\u201c\u51e1\u4e8b\u90fd\u8981\u89c4\u89c4\u77e9\u77e9\u5730\u6309\u7740\u6b21\u5e8f\u884c\u201d\uff08\u6797\u524d14:40\uff09',
    '\u5b9e\u8df5\u6848\u4f8b\uff1a\u6d88\u5c71\u6559\u4f1a\u5bf9\u5916\u6d3e\u8bb2\u5458\u8fdb\u884c\u795e\u5b66\u4e0e\u6cd5\u89c4\u53cc\u91cd\u8003\u6838'
]
add_bullets(s, Inches(0.8), Inches(2.5), Inches(11.733), Inches(4.5), items, size=18)
add_page_num(s, 3)

# ==================== Slide 4: Place ====================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, WHITE)
add_rect(s, Inches(0), Inches(0), Inches(13.333), Inches(0.08), DARK_BLUE)
add_rect(s, Inches(0), Inches(7.3), Inches(13.333), Inches(0.2), LIGHT_BLUE)

add_text(s, Inches(0.8), Inches(0.5), Inches(11.733), Inches(0.8),
         '\u4e8c\u3001\u4e25\u5b88\u201c\u5730\u201d\u7684\u8fb9\u754c', size=32, color=DARK_BLUE, bold=True)
add_rect(s, Inches(0.8), Inches(1.2), Inches(2), Inches(0.06), GOLD)
add_text(s, Inches(0.8), Inches(1.6), Inches(11.733), Inches(0.6),
         '\u843d\u5b9e\u767b\u8bb0\u5236\uff0c\u786e\u4fdd\u573a\u6240\u5408\u89c4', size=22, color=DARK_BLUE, bold=True)

items = [
    '\u96c6\u4f53\u5b97\u6559\u6d3b\u52a8\u987b\u5728\u4f9d\u6cd5\u7b79\u529e\u8bbe\u7acb\u548c\u767b\u8bb0\u7684\u573a\u6240\u5185\u5f00\u5c55',
    '\u4e25\u7981\u5728\u4f4f\u5b85\u3001\u5546\u94fa\u7b49\u975e\u5b97\u6559\u6d3b\u52a8\u573a\u6240\u79c1\u8bbe\u805a\u4f1a\u70b9',
    '\u786e\u6709\u9700\u8981\u7684\uff0c\u53ef\u7533\u8bf7\u6307\u5b9a\u5b97\u6559\u4e34\u65f6\u6d3b\u52a8\u5730\u70b9',
    '\u79c1\u8bbe\u70b9\u5b89\u5168\u9690\u60a3\u7a81\u51fa\uff0c\u6613\u53d7\u5f02\u7aef\u6e17\u900f',
    '\u5b9e\u8df5\u63aa\u65bd\uff1a\u6c11\u5b97\u3001\u6d88\u9632\u90e8\u95e8\u5e38\u6001\u5316\u8054\u5408\u5de1\u67e5'
]
add_bullets(s, Inches(0.8), Inches(2.5), Inches(11.733), Inches(4.5), items, size=18)
add_page_num(s, 4)

# ==================== Slide 5: Doctrine ====================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, WHITE)
add_rect(s, Inches(0), Inches(0), Inches(13.333), Inches(0.08), DARK_BLUE)
add_rect(s, Inches(0), Inches(7.3), Inches(13.333), Inches(0.2), LIGHT_BLUE)

add_text(s, Inches(0.8), Inches(0.5), Inches(11.733), Inches(0.8),
         '\u4e09\u3001\u628a\u52b3\u201c\u9053\u201d\u7684\u65b9\u5411', size=32, color=DARK_BLUE, bold=True)
add_rect(s, Inches(0.8), Inches(1.2), Inches(2), Inches(0.06), GOLD)
add_text(s, Inches(0.8), Inches(1.6), Inches(11.733), Inches(0.6),
         '\u843d\u5b9e\u4e2d\u56fd\u5316\uff0c\u786e\u4fdd\u4fe1\u4ef0\u6b63\u5411', size=22, color=DARK_BLUE, bold=True)

items = [
    '\u575a\u6301\u4e2d\u56fd\u5316\u65b9\u5411 \u2260 \u66f4\u6539\u6838\u5fc3\u771f\u7406',
    '\u4ee5\u793e\u4f1a\u4e3b\u4e49\u6838\u5fc3\u4ef7\u503c\u89c2\u4e3a\u5f15\u9886\uff0c\u589e\u8fdb\u201c\u4e94\u4e2a\u8ba4\u540c\u201d',
    '\u6df1\u5165\u6316\u6398\u6559\u4e49\u4e2d\u6709\u5229\u4e8e\u793e\u4f1a\u548c\u8c10\u3001\u65f6\u4ee3\u8fdb\u6b65\u7684\u5185\u5bb9',
    '\u201c\u5411\u4ec0\u4e48\u6837\u7684\u4eba\uff0c\u6211\u5c31\u4f5c\u4ec0\u4e48\u6837\u7684\u4eba\u201d\uff08\u6797\u524d9:22\uff09',
    '\u5b9e\u8df5\u54c1\u724c\uff1a\u201c\u4e4b\u6c5f\u95ee\u9053\u201d\u2014\u2014\u57fa\u7763\u6559\u4e2d\u56fd\u5316\u6d59\u6c5f\u63a2\u7d22'
]
add_bullets(s, Inches(0.8), Inches(2.5), Inches(11.733), Inches(4.5), items, size=18)
add_page_num(s, 5)

# ==================== Slide 6: Doctrine (cont.) ====================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, WHITE)
add_rect(s, Inches(0), Inches(0), Inches(13.333), Inches(0.08), DARK_BLUE)
add_rect(s, Inches(0), Inches(7.3), Inches(13.333), Inches(0.2), LIGHT_BLUE)

add_text(s, Inches(0.8), Inches(0.5), Inches(11.733), Inches(0.8),
         '\u4e09\u3001\u628a\u52b3\u201c\u9053\u201d\u7684\u65b9\u5411\uff08\u7eed\uff09', size=32, color=DARK_BLUE, bold=True)
add_rect(s, Inches(0.8), Inches(1.2), Inches(2), Inches(0.06), GOLD)
add_text(s, Inches(0.8), Inches(1.6), Inches(11.733), Inches(0.6),
         '\u8ba9\u6c38\u6052\u771f\u7406\u62d9\u6839\u4e2d\u534e\u3001\u63a5\u6d59\u6c5f\u5730\u6c14', size=22, color=DARK_BLUE, bold=True)

items = [
    '\u4e3e\u529e\u795e\u5b66\u601d\u60f3\u7814\u8ba8\u4f1a',
    '\u63a8\u52a8\u4e2d\u534e\u4f18\u79c0\u4f20\u7edf\u6587\u5316\u8fdb\u573a\u6240',
    '\u676d\u5dde\u57fa\u7763\u6559\u4eb2\u60c5\u8001\u4eba\u670d\u52a1\u4e2d\u5fc3\u3001\u601d\u6f84\u5802\u6301\u7eed20\u5e74\u4e66\u753b\u4e49\u5356\u670d\u52a1\u793e\u4f1a',
    '\u7528\u672c\u571f\u6545\u4e8b\u3001\u65f6\u4ee3\u8bdd\u8bed\u4f20\u8bb2\u4fe1\u4ef0\uff0c\u66f4\u8d34\u5408\u793e\u4f1a\u3001\u8d34\u8fd1\u4eba\u5fc3',
    '\u5b9e\u73b0\u4fe1\u4ef0\u7eaf\u6b63\u4e0e\u793e\u4f1a\u548c\u8c10\u7684\u7f8e\u7f8e\u4e0e\u5171'
]
add_bullets(s, Inches(0.8), Inches(2.5), Inches(11.733), Inches(4.5), items, size=18)
add_page_num(s, 6)

# ==================== Slide 7: Activities ====================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, WHITE)
add_rect(s, Inches(0), Inches(0), Inches(13.333), Inches(0.08), DARK_BLUE)
add_rect(s, Inches(0), Inches(7.3), Inches(13.333), Inches(0.2), LIGHT_BLUE)

add_text(s, Inches(0.8), Inches(0.5), Inches(11.733), Inches(0.8),
         '\u56db\u3001\u89c4\u8303\u201c\u4e8b\u201d\u7684\u7a0b\u5e8f', size=32, color=DARK_BLUE, bold=True)
add_rect(s, Inches(0.8), Inches(1.2), Inches(2), Inches(0.06), GOLD)
add_text(s, Inches(0.8), Inches(1.6), Inches(11.733), Inches(0.6),
         '\u843d\u5b9e\u5ba1\u6279\u5236\uff0c\u786e\u4fdd\u6d3b\u52a8\u5b89\u5168', size=22, color=DARK_BLUE, bold=True)

items = [
    '\u201c\u8981\u8131\u79bb\u4e0d\u6309\u89c4\u77e9\u800c\u884c\u7684\u4eba\u201d\uff08\u8d34\u540e3:6\uff09',
    '\u5927\u578b\u5b97\u6559\u6d3b\u52a8\u987b\u63d0\u524d30\u65e5\u62a5\u6279',
    '\u5236\u5b9a\u8be6\u6848\u9884\u6848\uff0c\u843d\u5b9e\u5b89\u5168\u8d23\u4efb',
    '\u505a\u597d\u4eba\u6d41\u63a7\u7ba1\u3001\u6d88\u9632\u5e94\u6025\u7b49\u4fdd\u969c',
    '\u843d\u5b9e\u201c\u6d3b\u52a8\u5b89\u5168\u4e3b\u8d23\u4eba\u201d\u5236\u5ea6\uff0c\u7ec6\u5316\u9632\u63a7\uff0c\u6392\u67e5\u98ce\u9669'
]
add_bullets(s, Inches(0.8), Inches(2.5), Inches(11.733), Inches(4.5), items, size=18)
add_page_num(s, 7)

# ==================== Slide 8: Law ====================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, WHITE)
add_rect(s, Inches(0), Inches(0), Inches(13.333), Inches(0.08), DARK_BLUE)
add_rect(s, Inches(0), Inches(7.3), Inches(13.333), Inches(0.2), LIGHT_BLUE)

add_text(s, Inches(0.8), Inches(0.5), Inches(11.733), Inches(0.8),
         '\u4e94\u3001\u660e\u6670\u201c\u6cd5\u201d\u7684\u7ea2\u7ebf', size=32, color=DARK_BLUE, bold=True)
add_rect(s, Inches(0.8), Inches(1.2), Inches(2), Inches(0.06), GOLD)
add_text(s, Inches(0.8), Inches(1.6), Inches(11.733), Inches(0.6),
         '\u843d\u5b9e\u6cd5\u6cbb\u5316\uff0c\u786e\u4fdd\u9886\u57df\u5b89\u5b9a', size=22, color=DARK_BLUE, bold=True)

items = [
    '\u300a\u5b97\u6559\u4e8b\u52a1\u6761\u4f8b\u300b\u5212\u5b9a\u201c\u5341\u4e2a\u7981\u6b62\u201d\u7ea2\u7ebf',
    '\u4e25\u7981\u6821\u56ed\u4f20\u6559\u3001\u672a\u6210\u5e74\u4eba\u53c2\u4e0e\u3001\u5883\u5916\u6e17\u900f\u3001\u64c5\u81ea\u8bbe\u70b9\u7b49',
    '\u4e09\u91cd\u4fdd\u969c\uff1a',
    '  \u00b7 \u5bf9\u9752\u5c11\u5e74\uff1a\u4fdd\u969c\u63a5\u53d7\u4e49\u52a1\u6559\u80b2\u3001\u5065\u5eb7\u6210\u957f\u7684\u6cd5\u6cbb\u5c4f\u969c',
    '  \u00b7 \u5bf9\u4fe1\u4f17\uff1a\u5e2e\u52a9\u660e\u8fa8\u662f\u975e\u3001\u62b5\u6321\u8bef\u5bfc\u3001\u6301\u5b88\u7eaf\u6b63\u7684\u575a\u56fa\u76fe\u724c',
    '  \u00b7 \u5bf9\u6559\u4f1a\uff1a\u575a\u6301\u72ec\u7acb\u81ea\u529e\u3001\u4e13\u6ce8\u7075\u6027\u7267\u517b\u3001\u5b9e\u73b0\u5065\u5eb7\u4f20\u627f\u7684\u6839\u672c\u51c6\u5219'
]
add_bullets(s, Inches(0.8), Inches(2.5), Inches(11.733), Inches(4.5), items, size=18)
add_page_num(s, 8)

# ==================== Slide 9: Conclusion ====================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, LIGHT_GRAY)
add_rect(s, Inches(0), Inches(0), Inches(13.333), Inches(0.08), DARK_BLUE)
add_rect(s, Inches(0), Inches(7.3), Inches(13.333), Inches(0.2), LIGHT_BLUE)

add_text(s, Inches(0.8), Inches(0.5), Inches(11.733), Inches(0.8),
         '\u7ed3\u8bed', size=32, color=DARK_BLUE, bold=True)
add_rect(s, Inches(0.8), Inches(1.2), Inches(2), Inches(0.06), GOLD)
add_text(s, Inches(0.8), Inches(1.6), Inches(11.733), Inches(0.6),
         '\u6cd5\u6cbb\u62a4\u822a\uff0c\u4fe1\u4ef0\u5728\u79e9\u5e8f\u4e2d\u81ea\u7531', size=22, color=DARK_BLUE, bold=True)

items = [
    '\u575a\u5b88\u7231\u56fd\u7231\u6559\u3001\u9075\u89c4\u5b88\u6cd5\u3001\u6b63\u4fe1\u6b63\u884c',
    '\u65e2\u6301\u5b88\u7eaf\u6b63\u771f\u9053\uff0c\u4e5f\u987a\u670d\u56fd\u5bb6\u6743\u957f\u6cd5\u5ea6',
    '\u8ba9\u4fe1\u4ef0\u5728\u6cd5\u6cbb\u9633\u5149\u4e0b\u5065\u5eb7\u4f20\u627f',
    '\u7acb\u8db3\u6d59\u6c5f\u3001\u6df1\u8015\u57fa\u5c42\u3001\u670d\u52a1\u793e\u4f1a',
    '\u79ef\u6781\u52a9\u529b\u5b97\u6559\u4e2d\u56fd\u5316\u6d59\u6c5f\u793a\u8303\u7701\u5efa\u8bbe',
    '\u4e3a\u6cd5\u6cbb\u6d59\u6c5f\u3001\u5e73\u5b89\u6d59\u6c5f\u5efa\u8bbe\u6301\u7eed\u8d21\u732e\u6b63\u5411\u529b\u91cf'
]
add_bullets(s, Inches(0.8), Inches(2.5), Inches(11.733), Inches(4.5), items, size=18)
add_page_num(s, 9)

# ==================== Slide 10: Thank You ====================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, DARK_BLUE)
add_rect(s, Inches(0), Inches(0), Inches(13.333), Inches(0.08), GOLD)
add_rect(s, Inches(4.5), Inches(3.0), Inches(4.333), Inches(0.04), GOLD)

add_text(s, Inches(1), Inches(3.3), Inches(11.333), Inches(1.5),
         '\u5ba3\u8bb2\u5b8c\u6bd5\n\u8c22\u8c22\u5927\u5bb6\uff01', size=44, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(s, Inches(1), Inches(5.5), Inches(11.333), Inches(0.6),
         '\u4f9d\u6cd5\u800c\u6cbb  \u5411\u5fc3\u800c\u884c', size=18, color=GOLD, align=PP_ALIGN.CENTER)
add_page_num(s, 10)

# Save
output = '/home/admin/.openclaw/workspace/ppt-outlines/基督教中国化宣讲.pptx'
prs.save(output)
print(f'PPT generated: {output}')
