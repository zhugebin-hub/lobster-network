#!/usr/bin/env python3
"""生成硕士学位论文答辩PPT"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ==================== 配置 ====================
OUTPUT_DIR = "/home/admin/.openclaw/workspace"
OUTPUT_FILE = os.path.join(OUTPUT_DIR, "基于博弈与CatBoost的分布式任务分配及网络切片资源预测_答辩PPT.pptx")

# 配色方案 - 学术蓝
COLOR_PRIMARY = RGBColor(0x1B, 0x3A, 0x7A)      # 深蓝
COLOR_SECONDARY = RGBColor(0x2E, 0x75, 0xB6)     # 中蓝
COLOR_ACCENT = RGBColor(0x5B, 0x9B, 0xD5)        # 浅蓝
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_DARK = RGBColor(0x33, 0x33, 0x33)
COLOR_GRAY = RGBColor(0x66, 0x66, 0x66)
COLOR_LIGHT_BG = RGBColor(0xE8, 0xF0, 0xFE)
COLOR_GOLD = RGBColor(0xC5, 0x96, 0x16)

# 字体
FONT_CN = "微软雅黑"
FONT_EN = "Calibri"

# ==================== 工具函数 ====================

def set_slide_bg(slide, color):
    """设置幻灯片背景色"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_rect(slide, left, top, width, height, fill_color, line_color=None):
    """添加矩形形状"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18, 
                 font_color=COLOR_DARK, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name=FONT_CN, line_spacing=1.3):
    """添加文本框"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(4)
    # 行距
    p.line_spacing = Pt(font_size * line_spacing)
    return txBox


def add_multiline_text(slide, left, top, width, height, lines, font_size=16,
                       font_color=COLOR_DARK, bold=False, alignment=PP_ALIGN.LEFT,
                       font_name=FONT_CN, bullet=False, line_spacing=1.5):
    """添加多行文本"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        if bullet:
            p.text = "  •  " + line
        else:
            p.text = line
        
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.bold = bold
        p.font.name = font_name
        p.alignment = alignment
        p.space_after = Pt(4)
        p.line_spacing = Pt(font_size * line_spacing)
    
    return txBox


def add_title_bar(slide, title_text, subtitle_text=""):
    """添加标题栏 - 深蓝背景 + 白色文字"""
    # 顶部深蓝条
    add_shape_rect(slide, Inches(0), Inches(0), Inches(10), Inches(1.4), COLOR_PRIMARY)
    # 底部装饰线
    add_shape_rect(slide, Inches(0), Inches(1.4), Inches(10), Inches(0.06), COLOR_GOLD)
    
    # 标题文字
    add_text_box(slide, Inches(0.5), Inches(0.25), Inches(9), Inches(0.7),
                 title_text, font_size=32, font_color=COLOR_WHITE, bold=True,
                 alignment=PP_ALIGN.LEFT)
    
    if subtitle_text:
        add_text_box(slide, Inches(0.5), Inches(0.85), Inches(9), Inches(0.5),
                     subtitle_text, font_size=16, font_color=RGBColor(0xBB, 0xD5, 0xF0),
                     bold=False, alignment=PP_ALIGN.LEFT)


def add_section_card(slide, left, top, width, height, title, content_lines, 
                     card_color=COLOR_LIGHT_BG, title_color=COLOR_PRIMARY):
    """添加内容卡片"""
    # 卡片背景
    shape = add_shape_rect(slide, left, top, width, height, card_color)
    shape.line.color.rgb = COLOR_ACCENT
    shape.line.width = Pt(1)
    
    # 左侧装饰条
    add_shape_rect(slide, left, top, Inches(0.08), height, COLOR_SECONDARY)
    
    # 标题
    add_text_box(slide, left + Inches(0.25), top + Inches(0.15), 
                 width - Inches(0.4), Inches(0.4),
                 title, font_size=18, font_color=title_color, bold=True)
    
    # 内容
    add_multiline_text(slide, left + Inches(0.25), top + Inches(0.55),
                       width - Inches(0.4), height - Inches(0.7),
                       content_lines, font_size=14, font_color=COLOR_DARK,
                       bullet=True, line_spacing=1.4)


def add_numbered_item(slide, left, top, width, number, title, content, num_color=COLOR_SECONDARY):
    """添加带编号的项目"""
    # 编号圆形
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top + Inches(0.05), 
                                    Inches(0.5), Inches(0.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = num_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = str(number)
    p.font.size = Pt(20)
    p.font.color.rgb = COLOR_WHITE
    p.font.bold = True
    p.font.name = FONT_CN
    p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    # 标题
    add_text_box(slide, left + Inches(0.65), top, width - Inches(0.65), Inches(0.35),
                 title, font_size=18, font_color=COLOR_PRIMARY, bold=True)
    
    # 内容
    add_multiline_text(slide, left + Inches(0.65), top + Inches(0.35), 
                       width - Inches(0.65), Inches(2.5),
                       content, font_size=14, font_color=COLOR_DARK,
                       bullet=True, line_spacing=1.4)


# ==================== 创建PPT ====================

print("开始创建PPT...")
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# ==================== 第1页：封面 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
set_slide_bg(slide, COLOR_WHITE)

# 顶部装饰条
add_shape_rect(slide, Inches(0), Inches(0), Inches(10), Inches(0.15), COLOR_PRIMARY)

# 学校/学院信息
add_text_box(slide, Inches(1), Inches(1.0), Inches(8), Inches(0.5),
             "浙江省硕士学位论文答辩", font_size=22, font_color=COLOR_SECONDARY, 
             bold=True, alignment=PP_ALIGN.CENTER)

# 主标题背景
add_shape_rect(slide, Inches(0.8), Inches(1.8), Inches(8.4), Inches(2.2), COLOR_PRIMARY)
add_shape_rect(slide, Inches(0.8), Inches(3.8), Inches(8.4), Inches(0.08), COLOR_GOLD)

# 主标题
add_text_box(slide, Inches(1.2), Inches(1.95), Inches(7.6), Inches(1.8),
             "基于博弈与CatBoost的\n分布式多智能体任务分配及网络切片资源预测研究",
             font_size=28, font_color=COLOR_WHITE, bold=True,
             alignment=PP_ALIGN.CENTER, line_spacing=1.4)

# 分隔线
add_shape_rect(slide, Inches(3), Inches(4.3), Inches(4), Inches(0.03), COLOR_ACCENT)

# 论文信息
info_lines = [
    "学科专业：信息与通信工程（081000）",
    "研究方向：计算机网络",
]
add_multiline_text(slide, Inches(2), Inches(4.6), Inches(6), Inches(1.0),
                   info_lines, font_size=16, font_color=COLOR_DARK,
                   alignment=PP_ALIGN.CENTER, line_spacing=1.8)

# 底部信息
add_text_box(slide, Inches(1), Inches(5.8), Inches(8), Inches(0.5),
             "答辩人：__________    导师：__________    日期：2026年6月",
             font_size=16, font_color=COLOR_GRAY, alignment=PP_ALIGN.CENTER)

# 底部装饰条
add_shape_rect(slide, Inches(0), Inches(7.35), Inches(10), Inches(0.15), COLOR_PRIMARY)


# ==================== 第2页：目录 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, COLOR_WHITE)
add_title_bar(slide, "目  录", "Contents")

toc_items = [
    ("01", "研究背景与意义", "Research Background and Significance"),
    ("02", "国内外研究现状", "Literature Review"),
    ("03", "研究内容与创新点", "Research Content and Contributions"),
    ("04", "理论架构", "Theoretical Framework"),
    ("05", "核心工作一：分布式任务分配算法", "Distributed Task Allocation Algorithm"),
    ("06", "核心工作二：网络切片资源预测", "Network Slicing Resource Prediction"),
    ("07", "实验结果与分析", "Experimental Results"),
    ("08", "总结与展望", "Conclusion and Future Work"),
]

for i, (num, title, en_title) in enumerate(toc_items):
    row = i // 2
    col = i % 2
    left = Inches(0.8) + col * Inches(4.8)
    top = Inches(1.8) + row * Inches(1.25)
    
    # 编号
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top + Inches(0.05),
                                    Inches(0.55), Inches(0.55))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_SECONDARY
    shape.line.fill.background()
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(18)
    p.font.color.rgb = COLOR_WHITE
    p.font.bold = True
    p.font.name = FONT_EN
    p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    # 标题
    add_text_box(slide, left + Inches(0.7), top, Inches(3.5), Inches(0.3),
                 title, font_size=17, font_color=COLOR_DARK, bold=True)
    add_text_box(slide, left + Inches(0.7), top + Inches(0.3), Inches(3.5), Inches(0.25),
                 en_title, font_size=11, font_color=COLOR_GRAY)


# ==================== 第3页：研究背景 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, COLOR_WHITE)
add_title_bar(slide, "研究背景与意义", "Research Background and Significance")

# 左侧卡片
add_section_card(slide, Inches(0.5), Inches(1.8), Inches(4.3), Inches(5.2),
                 "多智能体系统背景", [
                     "分布式、并行化、自适应性",
                     "广泛应用于机器人协作、无人机编队、智能交通",
                     "任务分配是NP-hard核心问题",
                     "传统集中式方法存在单点故障风险",
                     "动态环境下需要自适应分布式决策",
                 ])

# 右侧卡片
add_section_card(slide, Inches(5.2), Inches(1.8), Inches(4.3), Inches(5.2),
                 "网络切片背景", [
                     "5G网络催生网络切片技术",
                     "eMBB / URLLC / mMTC 多业务并行",
                     "切片性能受时延、丢包、吞吐量波动影响",
                     "准确预测切片状态是智能资源管理关键",
                     "SDN/NFV 提供灵活可定制的网络能力",
                 ])


# ==================== 第4页：研究问题 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, COLOR_WHITE)
add_title_bar(slide, "研究问题与挑战", "Research Problems and Challenges")

# 两个核心问题
add_numbered_item(slide, Inches(0.8), Inches(1.8), Inches(8.5), 1,
                  "问题一：动态环境下的分布式任务分配", [
                      "任务动态引入，奖励机制动态变化",
                      "通信受限，信息不完全",
                      "传统方法易陷入局部最优",
                      "需要兼顾效用最大化与通信开销最小化",
                  ])

add_numbered_item(slide, Inches(0.8), Inches(4.2), Inches(8.5), 2,
                  "问题二：网络切片资源需求预测", [
                      "资源需求波动大，业务类型多样",
                      "多切片并发运行，负载动态变化",
                      "预测精度与实时性难以兼顾",
                      "需要为后续资源调度提供可靠数据支撑",
                  ])


# ==================== 第5页：国内外研究现状 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, COLOR_WHITE)
add_title_bar(slide, "国内外研究现状", "Literature Review")

# 左侧：任务分配
add_section_card(slide, Inches(0.5), Inches(1.8), Inches(4.3), Inches(5.2),
                 "分布式多智能体任务分配", [
                     "群体智能方法：PSO、遗传算法（集中式局限）",
                     "基于市场机制：拍卖、竞标（通信开销大）",
                     "基于一致性算法：高通信负担",
                     "博弈论方法：效率高、可扩展性好",
                     "联盟形成博弈：自然描述多智能体协作",
                     "现有方法在动态环境下适应性有限",
                 ])

# 右侧：网络切片
add_section_card(slide, Inches(5.2), Inches(1.8), Inches(4.3), Inches(5.2),
                 "网络切片资源预测", [
                     "MILP模型：计算复杂度高，扩展性有限",
                     "GAN/LSTM：GANSlicing框架，链路资源预测",
                     "Bi-LSTM+Attention：VNF负载预测",
                     "GCN-GRU：物联网场景链路特征预测",
                     "深度强化学习：A2C、DQN资源分配",
                     "多数方法对多切片动态场景适应性不足",
                 ])


# ==================== 第6页：研究内容与创新点 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, COLOR_WHITE)
add_title_bar(slide, "研究内容与创新点", "Research Content and Contributions")

# 创新点1
add_shape_rect(slide, Inches(0.5), Inches(1.8), Inches(9), Inches(2.5), COLOR_LIGHT_BG)
add_shape_rect(slide, Inches(0.5), Inches(1.8), Inches(0.08), Inches(2.5), COLOR_SECONDARY)

add_text_box(slide, Inches(0.8), Inches(1.9), Inches(8.5), Inches(0.4),
             "创新点一：基于匿名享乐博弈的分布式多智能体任务分配算法（I_GRAPE）",
             font_size=18, font_color=COLOR_PRIMARY, bold=True)

add_multiline_text(slide, Inches(0.8), Inches(2.4), Inches(8.5), Inches(1.8), [
    "• 将任务分配建模为匿名享乐博弈，智能体仅依据任务特征及联盟规模决策",
    "• 引入匈牙利算法实现全局最优初始化，Q-Learning自适应调整多目标权重",
    "• 融合对数线性学习策略，平衡探索与利用，保证收敛至纳什均衡",
    "• 避免依赖智能体身份信息，有效降低通信开销",
], font_size=14, line_spacing=1.5)

# 创新点2
add_shape_rect(slide, Inches(0.5), Inches(4.6), Inches(9), Inches(2.5), COLOR_LIGHT_BG)
add_shape_rect(slide, Inches(0.5), Inches(4.6), Inches(0.08), Inches(2.5), COLOR_SECONDARY)

add_text_box(slide, Inches(0.8), Inches(4.7), Inches(8.5), Inches(0.4),
             "创新点二：基于Optuna超参数优化的CatBoost网络切片资源预测方法",
             font_size=18, font_color=COLOR_PRIMARY, bold=True)

add_multiline_text(slide, Inches(0.8), Inches(5.2), Inches(8.5), Inches(1.8), [
    "• 利用CatBoost有序提升机制和对称树结构，缓解梯度偏差与预测偏移",
    "• 引入Optuna框架实现超参数自动化搜索，提升模型泛化能力",
    "• 构建时间一致性约束机制，符合实际网络部署条件",
    "• 基于规则驱动的影响等级自动标注，实现高语义标签化",
], font_size=14, line_spacing=1.5)


# ==================== 第7页：理论架构 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, COLOR_WHITE)
add_title_bar(slide, "基于博弈的SDN驱动智能网络架构", 
              "Game Theory-based SDN-driven Intelligent Network Architecture")

# 三层架构示意
layers = [
    ("应用层", COLOR_PRIMARY, [
        "网络切片管理",
        "虚拟化服务",
        "安全管理",
        "元业务",
    ]),
    ("控制层（SDN智能体）", COLOR_SECONDARY, [
        "网络状态感知",
        "智能网络优化",
        "策略控制",
        "博弈决策",
    ]),
    ("转发层", COLOR_ACCENT, [
        "数据转发",
        "任务执行",
        "OpenFlow/NETCONF",
        "P4Runtime",
    ]),
]

for i, (name, color, items) in enumerate(layers):
    top = Inches(1.8) + i * Inches(1.6)
    
    # 层标题
    shape = add_shape_rect(slide, Inches(0.8), top, Inches(2.5), Inches(0.5), color)
    add_text_box(slide, Inches(0.8), top + Inches(0.05), Inches(2.5), Inches(0.4),
                 name, font_size=16, font_color=COLOR_WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER)
    
    # 功能模块
    for j, item in enumerate(items):
        left = Inches(3.6) + j * Inches(1.8)
        shape = add_shape_rect(slide, left, top + Inches(0.05), Inches(1.6), Inches(0.4), 
                               COLOR_LIGHT_BG)
        shape.line.color.rgb = color
        shape.line.width = Pt(1)
        add_text_box(slide, left, top + Inches(0.1), Inches(1.6), Inches(0.3),
                     item, font_size=12, font_color=COLOR_DARK, alignment=PP_ALIGN.CENTER)

# 接口说明
add_text_box(slide, Inches(0.8), Inches(6.8), Inches(8.5), Inches(0.4),
             "标准化REST接口  ←→  各层协同交互  ←→  实时数据驱动自适应资源优化",
             font_size=13, font_color=COLOR_GRAY, alignment=PP_ALIGN.CENTER)


# ==================== 第8页：匿名享乐博弈建模 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, COLOR_WHITE)
add_title_bar(slide, "匿名享乐博弈建模", "Anonymous Hedonic Game Modeling")

# 核心概念
add_section_card(slide, Inches(0.5), Inches(1.8), Inches(4.3), Inches(2.3),
                 "匿名享乐博弈特性", [
                     "效用仅取决于联盟规模，不依赖成员身份",
                     "联盟形成基于效用函数偏好排序",
                     "智能体通过协作最大化自身效用",
                     "收敛至纳什均衡：无人可单方面改善",
                 ])

# 效用函数
add_section_card(slide, Inches(5.2), Inches(1.8), Inches(4.3), Inches(2.3),
                 "效用函数设计", [
                     "Uij = 任务优先级 × 工作量 / 协作规模",
                     "− 时间成本权重 × 时间惩罚",
                     "− 距离成本权重 × 移动距离惩罚",
                     "+ 协作收益权重 × 协作效益",
                 ])

# 优化目标
add_section_card(slide, Inches(0.5), Inches(4.4), Inches(9), Inches(2.8),
                 "优化目标与约束", [
                     "最大化全局效用：所有智能体个体效用之和",
                     "约束1：一个智能体一次最多加入一个联盟",
                     "约束2：每个任务只能被一个联盟执行",
                     "约束3：任务一旦开始执行不可中途放弃",
                     "约束4：总执行时间不超过 Tmax",
                     "证明：该博弈构成序数潜在博弈，保证有限步收敛",
                 ])


# ==================== 第9页：I_GRAPE算法 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, COLOR_WHITE)
add_title_bar(slide, "I_GRAPE 算法实现", 
              "Improved Game-theoretic Resource Allocation for Multi-Agent Task Allocation")

# 算法流程步骤
steps = [
    ("① 初始化", "匈牙利算法全局最优\n初始任务划分", COLOR_PRIMARY),
    ("② Q-Learning\n权重调整", "状态：平均距离+\n完成进度", COLOR_SECONDARY),
    ("③ 对数线性\n学习策略", "概率选择任务\n自适应探索参数", COLOR_ACCENT),
    ("④ 联盟更新", "移除原联盟\n加入新联盟", RGBColor(0x70, 0xAD, 0x47)),
    ("⑤ 纳什均衡\n判定", "检查是否达到\n稳定状态", RGBColor(0xED, 0x7D, 0x31)),
    ("⑥ 广播执行", "广播联盟信息\n执行任务并更新", RGBColor(0x70, 0x30, 0xB0)),
]

for i, (title, desc, color) in enumerate(steps):
    col = i % 3
    row = i // 3
    left = Inches(0.5) + col * Inches(3.2)
    top = Inches(1.8) + row * Inches(2.5)
    
    # 步骤卡片
    shape = add_shape_rect(slide, left, top, Inches(2.9), Inches(2.1), COLOR_WHITE)
    shape.line.color.rgb = color
    shape.line.width = Pt(2)
    
    # 标题
    add_shape_rect(slide, left, top, Inches(2.9), Inches(0.5), color)
    add_text_box(slide, left, top + Inches(0.05), Inches(2.9), Inches(0.4),
                 title, font_size=15, font_color=COLOR_WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER)
    
    # 描述
    add_text_box(slide, left + Inches(0.15), top + Inches(0.6), Inches(2.6), Inches(1.4),
                 desc, font_size=13, font_color=COLOR_DARK, alignment=PP_ALIGN.CENTER,
                 line_spacing=1.5)
    
    # 箭头（非最后一个）
    if i < 5 and col < 2:
        add_text_box(slide, left + Inches(2.9), top + Inches(0.7), Inches(0.3), Inches(0.5),
                     "→", font_size=24, font_color=COLOR_GRAY, alignment=PP_ALIGN.CENTER)
    elif i < 5 and col == 2:
        add_text_box(slide, left + Inches(1.0), top + Inches(2.1), Inches(0.9), Inches(0.4),
                     "↓", font_size=24, font_color=COLOR_GRAY, alignment=PP_ALIGN.CENTER)


# ==================== 第10页：Q-Learning权重自适应 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, COLOR_WHITE)
add_title_bar(slide, "Q-Learning 多目标权重自适应机制",
              "Multi-objective Weight Adaptive Mechanism")

# 状态空间
add_section_card(slide, Inches(0.5), Inches(1.8), Inches(3.0), Inches(2.5),
                 "状态空间 S", [
                     "平均距离 d̄",
                     "平均完成进度 p̄",
                     "连续值离散化处理",
                     "简化状态空间维度",
                 ])

# 动作空间
add_section_card(slide, Inches(3.8), Inches(1.8), Inches(3.0), Inches(2.5),
                 "动作空间 A", [
                     "时间优先 (1.0,0.5,0.1)",
                     "距离优先 (0.5,1.0,0.1)",
                     "均衡策略 (0.7,0.7,0.7)",
                     "成本优先 (1.0,1.0,0.0)",
                     "协作优先 (0.5,0.5,1.0)",
                 ])

# 奖励函数
add_section_card(slide, Inches(7.1), Inches(1.8), Inches(2.4), Inches(2.5),
                 "奖励函数 R", [
                     "任务完成时间",
                     "总移动距离",
                     "协作效益",
                     "综合性能指标",
                 ])

# Q值更新公式
add_section_card(slide, Inches(0.5), Inches(4.6), Inches(9), Inches(2.5),
                 "Q值更新公式", [
                     "Q(s,a) ← Q(s,a) + η[r + γ·maxQ(s',a') − Q(s,a)]",
                     "η：学习速率（最优约0.1）",
                     "γ：折扣因子 [0,1]",
                     "r：即时奖励（综合性能指标）",
                     "无模型学习，通过试错自动优化权重组合",
                 ])


# ==================== 第11页：网络切片预测模型 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, COLOR_WHITE)
add_title_bar(slide, "基于CatBoost的网络切片资源预测",
              "Network Slicing Resource Prediction based on CatBoost")

# 问题定义
add_section_card(slide, Inches(0.5), Inches(1.8), Inches(4.3), Inches(2.3),
                 "预测问题定义", [
                     "输入：25维网络运行指标特征",
                     "输出：下一时间窗口影响等级",
                     "四等级：Adequate/Warning/Severe/Critical",
                     "时间一致性约束，避免信息泄露",
                 ])

# 标签生成
add_section_card(slide, Inches(5.2), Inches(1.8), Inches(4.3), Inches(2.3),
                 "规则驱动标签生成", [
                     "基于QoS指标：时延/丢包/吞吐量",
                     "引入应用权重（安全/效率/娱乐/通用）",
                     "最大影响原则：最严格QoS需求优先",
                     "无需人工标注，自动化标签生成",
                 ])

# CatBoost优势
add_section_card(slide, Inches(0.5), Inches(4.4), Inches(4.3), Inches(2.8),
                 "CatBoost 核心优势", [
                     "有序目标统计：高效处理类别特征",
                     "排序提升机制：缓解梯度偏差与预测偏移",
                     "对称决策树：降低复杂度，抑制过拟合",
                     "GPU加速支持：大规模数据快速建模",
                 ])

# Optuna优化
add_section_card(slide, Inches(5.2), Inches(4.4), Inches(4.3), Inches(2.8),
                 "Optuna 超参数优化", [
                     "基于贝叶斯优化的TPE算法",
                     "自适应搜索空间动态调整",
                     "提前终止低收益区域",
                     "自动搜索最优参数组合",
                 ])


# ==================== 第12页：实验设置 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, COLOR_WHITE)
add_title_bar(slide, "实验设置", "Experimental Setup")

# 任务分配实验
add_section_card(slide, Inches(0.5), Inches(1.8), Inches(4.3), Inches(5.0),
                 "任务分配实验", [
                     "仿真平台：基于Python的SPACE模拟平台",
                     "智能体规模：50（可扩展性测试10-100）",
                     "任务规模：20/30/40（可扩展性测试100-1000）",
                     "动态任务：每隔1000秒创建5个额外任务",
                     "对比算法：GRAPE / Greedy / CBBA",
                     "评估指标：效用/移动距离/完成时间",
                     "消融实验：N-QWAM（无权重调整）/ N-LLL（无对数线性学习）",
                     "统计检验：160次仿真，t检验+Cohen's d效应量",
                 ])

# 预测实验
add_section_card(slide, Inches(5.2), Inches(1.8), Inches(4.3), Inches(5.0),
                 "资源预测实验", [
                     "数据集：Kaggle网络服务质量数据（2万条）",
                     "特征：25维QoS指标（时延/丢包/吞吐量）",
                     "数据集划分：80%训练 / 20%测试（时间顺序）",
                     "对比模型：RandomForest / TabNet",
                     "评估指标：Accuracy / Precision / Recall",
                     "           F1-Macro / F1 Score（各等级）",
                     "           计算开销对比",
                     "超参数搜索：Optuna自动化调优",
                 ])


# ==================== 第13页：任务分配实验结果 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, COLOR_WHITE)
add_title_bar(slide, "任务分配实验结果", "Task Allocation Experimental Results")

# 消融实验
add_section_card(slide, Inches(0.5), Inches(1.8), Inches(4.3), Inches(2.5),
                 "消融实验结果", [
                     "完整I_GRAPE优于N-QWAM和N-LLL",
                     "任务规模20：提升1.08%（vs N-QWAM）",
                     "任务规模30：提升17.88%（vs N-QWAM）",
                     "任务规模40：提升10.75%（vs N-QWAM）",
                     "两大模块互补增益，协同作用显著",
                 ])

# 对比实验
add_section_card(slide, Inches(5.2), Inches(1.8), Inches(4.3), Inches(2.5),
                 "算法对比实验", [
                     "I_GRAPE vs GRAPE：提升7.47%-12.37%",
                     "I_GRAPE vs Greedy：提升5.85%-14.14%",
                     "I_GRAPE vs CBBA：提升8.90%-19.14%",
                     "移动距离最低，方差最小",
                     "任务完成速度最快，工作量下降最迅速",
                 ])

# 可扩展性
add_section_card(slide, Inches(0.5), Inches(4.6), Inches(9), Inches(2.6),
                 "可扩展性分析", [
                     "智能体10→100：平均移动距离400→40，良好规模扩展趋势",
                     "不同任务规模下移动距离分布稳定，鲁棒性强",
                     "t检验：所有规模组合均达统计极显著水平（p<0.001）",
                     "效应量分析：d=1.92-3.84，大至极大效应量",
                     "充分验证算法在大规模分布式多智能体系统中的潜力",
                 ])


# ==================== 第14页：资源预测实验结果 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, COLOR_WHITE)
add_title_bar(slide, "资源预测实验结果", "Resource Prediction Experimental Results")

# 主要结果
add_section_card(slide, Inches(0.5), Inches(1.8), Inches(4.3), Inches(2.5),
                 "整体预测性能", [
                     "CatBoost准确率：95.27%",
                     "F1-Macro：0.9524",
                     "vs RandomForest：准确率提升约0.73%",
                     "vs TabNet：准确率提升约11.4%",
                     "交叉验证与测试集差距小，无明显过拟合",
                 ])

# 各等级F1
add_section_card(slide, Inches(5.2), Inches(1.8), Inches(4.3), Inches(2.5),
                 "各影响等级F1分数", [
                     "Adequate等级：F1 > 0.95",
                     "Warning等级：F1 > 0.93",
                     "Severe等级：F1 > 0.95",
                     "Critical等级：F1 = 0.959",
                     "关键等级识别能力强，风险预警可靠",
                 ])

# 计算开销
add_section_card(slide, Inches(0.5), Inches(4.6), Inches(4.3), Inches(2.6),
                 "计算开销对比", [
                     "CatBoost：总运行时间较短",
                     "计算成本主要在超参数优化阶段",
                     "RandomForest：运行时间最短",
                     "TabNet：总运行时间最长，开销显著",
                     "CatBoost综合性能最优",
                 ])

# 关键结论
add_section_card(slide, Inches(5.2), Inches(4.6), Inches(4.3), Inches(2.6),
                 "关键结论", [
                     "CatBoost在结构化数据任务中优势明显",
                     "Optuna超参数调优策略有效",
                     "时间一致性约束机制保证模型可信性",
                     "深度学习模型（TabNet）在小规模数据上泛化不足",
                     "为网络资源调度和风险预警提供可靠决策支持",
                 ])


# ==================== 第15页：总结 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, COLOR_WHITE)
add_title_bar(slide, "研究总结", "Research Summary")

# 三大成果
add_numbered_item(slide, Inches(0.8), Inches(1.8), Inches(8.5), 1,
                  "SDN驱动智能网络架构", [
                      "融合博弈理论的三层架构（应用层/控制层/转发层）",
                      "为分布式资源管理与优化决策提供统一系统支撑",
                  ])

add_numbered_item(slide, Inches(0.8), Inches(3.8), Inches(8.5), 2,
                  "分布式多智能体任务分配算法", [
                      "匿名享乐博弈框架 + 匈牙利初始化 + Q-Learning权重自适应",
                      "对数线性学习策略保证收敛至纳什均衡",
                      "相比经典算法性能提升 5.85%-19.14%",
                  ])

add_numbered_item(slide, Inches(0.8), Inches(5.8), Inches(8.5), 3,
                  "网络切片资源预测方法", [
                      "CatBoost + Optuna超参数优化",
                      "时间一致性约束 + 规则驱动标签生成",
                      "准确率 95.27%，F1-Macro 0.9524",
                  ])


# ==================== 第16页：展望 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, COLOR_WHITE)
add_title_bar(slide, "未来工作展望", "Future Work")

# 四个方向
directions = [
    ("实时动态性", "考虑任务到达频率不稳定、网络状态快速波动等复杂实时动态场景，引入在线学习和自适应更新机制", COLOR_PRIMARY),
    ("多算法融合", "探索启发式优化与学习型算法相结合，建立多种智能算法之间的协同融合机制", COLOR_SECONDARY),
    ("多性能评估", "引入能耗开销、负载均衡性、系统公平性、QoS保障等多维度综合评估模式", COLOR_ACCENT),
    ("实际场景应用", "结合真实网络环境或实验测试平台验证，探索云数据中心管理、网络切片调度等实际应用场景", RGBColor(0x70, 0xAD, 0x47)),
]

for i, (title, desc, color) in enumerate(directions):
    row = i // 2
    col = i % 2
    left = Inches(0.5) + col * Inches(4.8)
    top = Inches(1.8) + row * Inches(2.5)
    
    # 卡片
    shape = add_shape_rect(slide, left, top, Inches(4.5), Inches(2.2), COLOR_WHITE)
    shape.line.color.rgb = color
    shape.line.width = Pt(2)
    
    # 顶部色条
    add_shape_rect(slide, left, top, Inches(4.5), Inches(0.5), color)
    
    # 标题
    add_text_box(slide, left, top + Inches(0.05), Inches(4.5), Inches(0.4),
                 title, font_size=18, font_color=COLOR_WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER)
    
    # 描述
    add_text_box(slide, left + Inches(0.2), top + Inches(0.65), Inches(4.1), Inches(1.4),
                 desc, font_size=14, font_color=COLOR_DARK, line_spacing=1.5)


# ==================== 第17页：致谢 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, COLOR_WHITE)

# 顶部装饰
add_shape_rect(slide, Inches(0), Inches(0), Inches(10), Inches(0.15), COLOR_PRIMARY)

# 致谢标题
add_shape_rect(slide, Inches(3), Inches(2.0), Inches(4), Inches(0.8), COLOR_PRIMARY)
add_text_box(slide, Inches(3), Inches(2.1), Inches(4), Inches(0.6),
             "感谢聆听！", font_size=36, font_color=COLOR_WHITE, bold=True,
             alignment=PP_ALIGN.CENTER)

# 装饰线
add_shape_rect(slide, Inches(3.5), Inches(3.2), Inches(3), Inches(0.03), COLOR_GOLD)

# 致谢内容
thanks_lines = [
    "感谢导师的悉心指导",
    "感谢评审专家的宝贵意见",
    "感谢实验室同门的帮助与支持",
]
add_multiline_text(slide, Inches(2.5), Inches(3.8), Inches(5), Inches(2.0),
                   thanks_lines, font_size=18, font_color=COLOR_DARK,
                   alignment=PP_ALIGN.CENTER, line_spacing=2.0)

# 底部装饰
add_shape_rect(slide, Inches(0), Inches(7.35), Inches(10), Inches(0.15), COLOR_PRIMARY)


# ==================== 保存 ====================
prs.save(OUTPUT_FILE)
print(f"PPT已保存至: {OUTPUT_FILE}")
print(f"总页数: {len(prs.slides)}")
