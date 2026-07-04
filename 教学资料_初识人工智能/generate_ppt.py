#!/usr/bin/env python3
"""生成《初识人工智能》PPT课件"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ===== 配色方案 =====
BLUE_DARK = RGBColor(0x1A, 0x23, 0x7E)
BLUE_MED = RGBColor(0x28, 0x3D, 0xA8)
BLUE_LIGHT = RGBColor(0x3F, 0x51, 0xB5)
ACCENT_ORANGE = RGBColor(0xFF, 0x6F, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
MED_GRAY = RGBColor(0x66, 0x66, 0x66)

def add_bg(slide, color):
    """设置幻灯片背景色"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape_bg(slide, left, top, width, height, color, alpha=None):
    """添加矩形背景"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if alpha:
        shape.fill.fore_color.brightness = alpha
    return shape

def set_text(tf, text, size=18, color=DARK_GRAY, bold=False, align=PP_ALIGN.LEFT, font_name="微软雅黑"):
    """设置文本框内容"""
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = font_name
    return p

def add_textbox(slide, left, top, width, height, text, size=18, color=DARK_GRAY, bold=False, align=PP_ALIGN.LEFT, font_name="微软雅黑"):
    """添加文本框"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    set_text(tf, text, size, color, bold, align, font_name)
    return txBox

def add_bullet_list(slide, left, top, width, height, items, size=16, color=DARK_GRAY, spacing=Pt(8)):
    """添加项目符号列表"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = spacing
        p.level = 0
        run = p.add_run()
        run.text = f"• {item}"
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.name = "微软雅黑"
    return txBox

def add_icon_text(slide, left, top, width, height, icon, title, desc, size=14):
    """添加图标+标题+描述的卡片"""
    # 卡片背景
    card = add_shape_bg(slide, left, top, width, height, WHITE)
    card.shadow.inherit = False

    # 图标
    add_textbox(slide, left + Inches(0.15), top + Inches(0.1), width, Inches(0.4), icon, size=22, align=PP_ALIGN.CENTER)
    # 标题
    add_textbox(slide, left + Inches(0.1), top + Inches(0.5), width - Inches(0.2), Inches(0.35), title, size=size+2, color=BLUE_DARK, bold=True, align=PP_ALIGN.CENTER)
    # 描述
    add_textbox(slide, left + Inches(0.15), top + Inches(0.85), width - Inches(0.3), height - Inches(1.0), desc, size=size, color=MED_GRAY, align=PP_ALIGN.CENTER)


# ============================================================
# 第1页：封面
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白版式
add_bg(slide, BLUE_DARK)

# 装饰线条
add_shape_bg(slide, Inches(0), Inches(2.5), Inches(13.333), Inches(0.04), ACCENT_ORANGE)

# 主标题
add_textbox(slide, Inches(1.5), Inches(2.8), Inches(10.3), Inches(1.5),
            "初识人工智能", size=52, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

# 副标题
add_textbox(slide, Inches(2), Inches(4.3), Inches(9.3), Inches(0.8),
            "信息技术 · 高一 · 人工智能初步", size=24, color=RGBColor(0xBB, 0xDE, 0xFB), align=PP_ALIGN.CENTER)

# 底部信息
add_textbox(slide, Inches(2), Inches(5.8), Inches(9.3), Inches(0.6),
            "授课教师：____________    日期：____________", size=16, color=RGBColor(0x90, 0xCA, 0xF9), align=PP_ALIGN.CENTER)

# 底部装饰线
add_shape_bg(slide, Inches(0), Inches(7.0), Inches(13.333), Inches(0.04), ACCENT_ORANGE)


# ============================================================
# 第2页：课堂导入
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

# 标题栏
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), BLUE_DARK)
add_textbox(slide, Inches(0.5), Inches(0.25), Inches(12), Inches(0.8),
            "猜一猜，这些是AI做的吗？", size=32, color=WHITE, bold=True)

# 说明文字
add_textbox(slide, Inches(1), Inches(1.8), Inches(11.3), Inches(0.6),
            "仔细观察下面的内容，你能分辨出哪些是人类创作的，哪些是AI生成的吗？",
            size=18, color=MED_GRAY, align=PP_ALIGN.CENTER)

# 4个展示框
boxes = [
    ("🎨", "AI绘画", "这幅风景画是AI用30秒生成的"),
    ("📰", "AI播报", "这段新闻是由AI语音合成的"),
    ("✍️", "AI写诗", "这首诗是AI根据关键词创作的"),
    ("🎵", "AI作曲", "这段音乐是AI独立创作的"),
]
for i, (icon, title, desc) in enumerate(boxes):
    x = Inches(1) + Inches(3) * i
    add_icon_text(slide, x, Inches(2.8), Inches(2.6), Inches(2.5), icon, title, desc)

# 底部问题
add_textbox(slide, Inches(1), Inches(6.0), Inches(11.3), Inches(0.8),
            "💡 你的判断：以上4个内容中，有 ____ 个是AI生成的。",
            size=20, color=BLUE_DARK, bold=True, align=PP_ALIGN.CENTER)


# ============================================================
# 第3页：学习目标
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), BLUE_DARK)
add_textbox(slide, Inches(0.5), Inches(0.25), Inches(12), Inches(0.8),
            "🎯 今天我们要学什么？", size=32, color=WHITE, bold=True)

goals = [
    ("🤔", "了解概念", "什么是人工智能？它有哪些特征？"),
    ("📜", "发展简史", "AI是如何诞生和发展的？经历了哪些阶段？"),
    ("🔍", "发现应用", "AI在我们的生活中有哪些应用？"),
    ("🛠️", "动手体验", "亲自体验AI工具，感受AI的能力"),
    ("💭", "思考影响", "AI给我们带来了哪些机遇与挑战？"),
]
for i, (icon, title, desc) in enumerate(goals):
    y = Inches(1.8) + Inches(1.0) * i
    # 序号圆
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.5), y + Inches(0.05), Inches(0.6), Inches(0.6))
    circle.fill.solid()
    circle.fill.fore_color.rgb = BLUE_MED
    circle.line.fill.background()
    set_text(circle.text_frame, f"{i+1}", size=20, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(2.3), y, Inches(2), Inches(0.5), icon + " " + title, size=20, color=BLUE_DARK, bold=True)
    add_textbox(slide, Inches(4.5), y + Inches(0.05), Inches(7), Inches(0.5), desc, size=16, color=MED_GRAY)


# ============================================================
# 第4页：什么是人工智能
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), BLUE_DARK)
add_textbox(slide, Inches(0.5), Inches(0.25), Inches(12), Inches(0.8),
            "什么是人工智能？", size=32, color=WHITE, bold=True)

# 通俗解释
add_shape_bg(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(2.0), RGBColor(0xE3, 0xF2, 0xFD))
add_textbox(slide, Inches(1.0), Inches(1.7), Inches(5.1), Inches(0.5), "💡 通俗说法", size=20, color=BLUE_DARK, bold=True)
add_textbox(slide, Inches(1.0), Inches(2.2), Inches(5.1), Inches(1.2),
            "让机器像人一样\n去感知、去学习、去思考、去行动",
            size=18, color=DARK_GRAY)

# 英文定义
add_shape_bg(slide, Inches(6.8), Inches(1.6), Inches(5.7), Inches(2.0), RGBColor(0xE3, 0xF2, 0xFD))
add_textbox(slide, Inches(7.0), Inches(1.7), Inches(5.3), Inches(0.5), "📖 英文定义", size=20, color=BLUE_DARK, bold=True)
add_textbox(slide, Inches(7.0), Inches(2.2), Inches(5.3), Inches(0.5),
            "Artificial Intelligence", size=16, color=BLUE_MED, bold=True)
add_textbox(slide, Inches(7.0), Inches(2.7), Inches(5.3), Inches(0.7),
            "简称 AI，是研究、开发用于模拟、\n延伸和扩展人类智能的理论、方法和技术。",
            size=14, color=MED_GRAY)

# 核心能力
add_textbox(slide, Inches(0.8), Inches(4.0), Inches(11.7), Inches(0.6),
            "🧠 人工智能的四大核心能力", size=22, color=BLUE_DARK, bold=True, align=PP_ALIGN.CENTER)

abilities = [
    ("👁️", "感知", "看、听、感知世界"),
    ("🧠", "理解", "读懂文字、听懂语言"),
    ("💡", "推理", "分析、判断、决策"),
    ("🤖", "行动", "执行任务、解决问题"),
]
for i, (icon, title, desc) in enumerate(abilities):
    x = Inches(1.0) + Inches(2.9) * i
    add_icon_text(slide, x, Inches(4.7), Inches(2.5), Inches(2.2), icon, title, desc, size=14)


# ============================================================
# 第5页：AI发展简史（时间轴）
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), BLUE_DARK)
add_textbox(slide, Inches(0.5), Inches(0.25), Inches(12), Inches(0.8),
            "📜 人工智能发展简史", size=32, color=WHITE, bold=True)

events = [
    ("1950", "图灵测试", "图灵提出机器智能判断标准", ACCENT_ORANGE),
    ("1956", "AI诞生", "达特茅斯会议，概念正式提出 ⭐", RGBColor(0x4C, 0xAF, 0x50)),
    ("1970s", "第一次低谷", "预期过高，技术受限", MED_GRAY),
    ("1980s", "专家系统", "专家系统兴起，AI复苏", RGBColor(0x03, 0xA1, 0x9E)),
    ("1990s", "第二次低谷", "专家系统局限性显现", MED_GRAY),
    ("2010s", "深度学习", "深度学习爆发，AI崛起 🚀", RGBColor(0xFF, 0x98, 0x00)),
    ("2020s", "大模型时代", "ChatGPT等引发AI全面爆发 🔥", ACCENT_ORANGE),
]

for i, (year, title, desc, color) in enumerate(events):
    y = Inches(1.6) + Inches(0.78) * i
    # 年份
    add_textbox(slide, Inches(0.8), y, Inches(1.2), Inches(0.5), year, size=16, color=color, bold=True)
    # 连接线
    if i < len(events) - 1:
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2.1), y + Inches(0.35), Inches(0.06), Inches(0.55))
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
        line.line.fill.background()
    # 圆点
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(2.0), y + Inches(0.12), Inches(0.35), Inches(0.35))
    dot.fill.solid()
    dot.fill.fore_color.rgb = color
    dot.line.fill.background()
    # 标题
    add_textbox(slide, Inches(2.5), y, Inches(3), Inches(0.45), title, size=17, color=DARK_GRAY, bold=True)
    # 描述
    add_textbox(slide, Inches(5.5), y + Inches(0.05), Inches(7), Inches(0.45), desc, size=14, color=MED_GRAY)


# ============================================================
# 第6页：AI的三个阶段
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), BLUE_DARK)
add_textbox(slide, Inches(0.5), Inches(0.25), Inches(12), Inches(0.8),
            "人工智能的三个阶段", size=32, color=WHITE, bold=True)

stages = [
    ("弱人工智能", "ANI", "专注某一特定领域\n如：AlphaGo、人脸识别\n语音助手、推荐算法", "✅\n已实现", BLUE_MED),
    ("强人工智能", "AGI", "具备人类级别智能\n能完成各种任务\n像人一样思考和行动", "❌\n未实现", RGBColor(0xFF, 0x98, 0x00)),
    ("超人工智能", "ASI", "在所有方面超越\n人类智能的AI\n存在于科幻作品中", "❌\n未实现", ACCENT_ORANGE),
]

for i, (title, abbr, desc, status, color) in enumerate(stages):
    x = Inches(0.8) + Inches(3.9) * i
    # 卡片
    card = add_shape_bg(slide, x, Inches(1.6), Inches(3.5), Inches(5.2), RGBColor(0xF8, 0xF9, 0xFA))
    card.corner_radius = Emu(10000)

    # 颜色条
    add_shape_bg(slide, x, Inches(1.6), Inches(3.5), Inches(0.08), color)

    # 阶段名称
    add_textbox(slide, x + Inches(0.2), Inches(1.9), Inches(3.1), Inches(0.5), title, size=24, color=color, bold=True, align=PP_ALIGN.CENTER)
    # 缩写
    add_textbox(slide, x + Inches(0.2), Inches(2.4), Inches(3.1), Inches(0.4), abbr, size=16, color=MED_GRAY, align=PP_ALIGN.CENTER)

    # 描述
    add_textbox(slide, x + Inches(0.3), Inches(3.0), Inches(2.9), Inches(2.5), desc, size=15, color=DARK_GRAY, align=PP_ALIGN.CENTER)

    # 状态标签
    add_shape_bg(slide, x + Inches(0.8), Inches(5.7), Inches(1.9), Inches(0.8), color)
    add_textbox(slide, x + Inches(0.8), Inches(5.7), Inches(1.9), Inches(0.8), status, size=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER)


# ============================================================
# 第7页：AI在身边——生活篇
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), BLUE_DARK)
add_textbox(slide, Inches(0.5), Inches(0.25), Inches(12), Inches(0.8),
            "🔍 AI在身边——生活篇", size=32, color=WHITE, bold=True)

apps = [
    ("📱", "手机解锁", "人脸识别\n指纹识别"),
    ("🎵", "个性推荐", "抖音/淘宝\n猜你喜欢"),
    ("🗣️", "语音助手", "Siri/小爱同学\n小度/天猫精灵"),
    ("📸", "智能相册", "自动分类\n人脸识别"),
    ("🏠", "智能家居", "扫地机器人\n智能音箱"),
]

for i, (icon, title, desc) in enumerate(apps):
    x = Inches(0.8) + Inches(2.4) * i
    add_icon_text(slide, x, Inches(1.6), Inches(2.1), Inches(2.8), icon, title, desc, size=14)

# 互动提示
add_shape_bg(slide, Inches(0.8), Inches(5.0), Inches(11.7), Inches(1.8), RGBColor(0xE3, 0xF2, 0xFD))
add_textbox(slide, Inches(1.2), Inches(5.1), Inches(11), Inches(0.5),
            "💬 课堂互动：你每天会用到哪些AI？", size=20, color=BLUE_DARK, bold=True)
add_textbox(slide, Inches(1.2), Inches(5.6), Inches(11), Inches(1.0),
            "从衣、食、住、行、学五个方面想一想，小组讨论并记录在学习任务单上。\n时间：5分钟",
            size=16, color=MED_GRAY)


# ============================================================
# 第8页：AI在身边——更多领域
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), BLUE_DARK)
add_textbox(slide, Inches(0.5), Inches(0.25), Inches(12), Inches(0.8),
            "🌍 AI无处不在", size=32, color=WHITE, bold=True)

fields = [
    ("🚗", "交通出行", "自动驾驶\n智能红绿灯\n路线规划"),
    ("🏥", "医疗健康", "AI辅助诊断\n药物研发\n医学影像分析"),
    ("🏭", "工业制造", "智能制造\n质量检测\n预测性维护"),
    ("🎮", "文化娱乐", "AI游戏NPC\n虚拟偶像\nAI配乐"),
    ("🛡️", "安全安防", "人脸识别\n行为分析\n智能监控"),
    ("📚", "教育学习", "智能批改\n个性化学习\nAI陪练"),
]

for i, (icon, title, desc) in enumerate(fields):
    col = i % 3
    row = i // 3
    x = Inches(0.8) + Inches(4.0) * col
    y = Inches(1.6) + Inches(2.5) * row
    add_icon_text(slide, x, y, Inches(3.5), Inches(2.2), icon, title, desc, size=13)


# ============================================================
# 第9页：AI工具体验
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), BLUE_DARK)
add_textbox(slide, Inches(0.5), Inches(0.25), Inches(12), Inches(0.8),
            "🛠️ AI工具体验站", size=32, color=WHITE, bold=True)

add_textbox(slide, Inches(1), Inches(1.4), Inches(11.3), Inches(0.5),
            "分组体验以下AI工具，完成任务单相关任务", size=18, color=MED_GRAY, align=PP_ALIGN.CENTER)

stations = [
    ("🎨", "AI创作站", "用AI生成一幅画", "文心一格\n通义万相", ACCENT_ORANGE),
    ("🗣️", "AI语音站", "和AI对话聊天", "小爱同学\n通义千问", BLUE_MED),
    ("🌐", "AI翻译站", "AI翻译对比体验", "通义千问\n文心一言", RGBColor(0x4C, 0xAF, 0x50)),
    ("👁️", "AI识别站", "AI看图识万物", "百度识图\n通义千问", RGBColor(0x9C, 0x27, 0xB0)),
]

for i, (icon, title, task, tools, color) in enumerate(stations):
    x = Inches(0.8) + Inches(3.0) * i
    card = add_shape_bg(slide, x, Inches(2.2), Inches(2.7), Inches(4.5), WHITE)

    # 顶部颜色条
    add_shape_bg(slide, x, Inches(2.2), Inches(2.7), Inches(0.06), color)

    add_textbox(slide, x + Inches(0.1), Inches(2.5), Inches(2.5), Inches(0.6), icon, size=40, align=PP_ALIGN.CENTER)
    add_textbox(slide, x + Inches(0.1), Inches(3.1), Inches(2.5), Inches(0.5), title, size=20, color=color, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, x + Inches(0.1), Inches(3.6), Inches(2.5), Inches(0.5), task, size=15, color=DARK_GRAY, align=PP_ALIGN.CENTER)

    # 工具推荐
    add_shape_bg(slide, x + Inches(0.3), Inches(4.3), Inches(2.1), Inches(1.2), RGBColor(0xF5, 0xF5, 0xF5))
    add_textbox(slide, x + Inches(0.3), Inches(4.35), Inches(2.1), Inches(0.3), "推荐工具", size=12, color=MED_GRAY, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, x + Inches(0.3), Inches(4.65), Inches(2.1), Inches(0.7), tools, size=14, color=DARK_GRAY, align=PP_ALIGN.CENTER)


# ============================================================
# 第10页：AI的优势与局限
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), BLUE_DARK)
add_textbox(slide, Inches(0.5), Inches(0.25), Inches(12), Inches(0.8),
            "AI很强，但不是万能的", size=32, color=WHITE, bold=True)

# 优势
add_shape_bg(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(5.2), RGBColor(0xE8, 0xF5, 0xE9))
add_textbox(slide, Inches(1.0), Inches(1.8), Inches(5.1), Inches(0.5), "✅ AI的优势", size=24, color=RGBColor(0x2E, 0x7D, 0x32), bold=True)

strengths = [
    "速度快，效率高",
    "处理大数据能力强",
    "不知疲倦，24小时工作",
    "精确度高，少犯错",
    "可重复执行，稳定性好",
    "能完成人类难以完成的任务",
]
add_bullet_list(slide, Inches(1.2), Inches(2.4), Inches(4.8), Inches(4.0), strengths, size=16, color=DARK_GRAY)

# 局限
add_shape_bg(slide, Inches(7.0), Inches(1.6), Inches(5.5), Inches(5.2), RGBColor(0xFF, 0xF3, 0xE0))
add_textbox(slide, Inches(7.2), Inches(1.8), Inches(5.1), Inches(0.5), "❌ AI的局限", size=24, color=ACCENT_ORANGE, bold=True)

limits = [
    "没有真正的情感和意识",
    "缺乏真正的创造力",
    "依赖数据质量和数量",
    "无法理解复杂的社会规则",
    "可能出现偏见和歧视",
    "无法为决策承担道德责任",
]
add_bullet_list(slide, Inches(7.4), Inches(2.4), Inches(4.8), Inches(4.0), limits, size=16, color=DARK_GRAY)


# ============================================================
# 第11页：思考与讨论
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), BLUE_DARK)
add_textbox(slide, Inches(0.5), Inches(0.25), Inches(12), Inches(0.8),
            "💭 AI时代，我们该怎么办？", size=32, color=WHITE, bold=True)

# 讨论话题
add_textbox(slide, Inches(1), Inches(1.6), Inches(11.3), Inches(0.5),
            "小组讨论（选择以下一个话题）：", size=20, color=BLUE_DARK, bold=True)

topics = [
    ("话题一", "AI会不会取代我们的工作？"),
    ("话题二", "AI绘画算不算真正的艺术创作？"),
    ("话题三", "我们应该害怕AI还是拥抱AI？"),
]

for i, (label, topic) in enumerate(topics):
    y = Inches(2.3) + Inches(1.2) * i
    add_shape_bg(slide, Inches(1.5), y, Inches(1.5), Inches(0.9), BLUE_MED)
    add_textbox(slide, Inches(1.5), y + Inches(0.15), Inches(1.5), Inches(0.6), label, size=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(3.2), y + Inches(0.15), Inches(8), Inches(0.6), topic, size=20, color=DARK_GRAY, bold=True)

# 引导要点
add_shape_bg(slide, Inches(1), Inches(5.8), Inches(11.3), Inches(1.4), RGBColor(0xE3, 0xF2, 0xFD))
add_textbox(slide, Inches(1.3), Inches(5.9), Inches(10.7), Inches(0.4), "💡 引导要点", size=18, color=BLUE_DARK, bold=True)
points = [
    "AI会取代部分重复性工作，但也会创造新的就业岗位",
    '关键是学会与AI协作，做AI的"主人"',
    "不是AI取代你，是先学会AI的人取代你",
]
add_bullet_list(slide, Inches(1.3), Inches(6.3), Inches(10.7), Inches(1.0), points, size=14, color=DARK_GRAY)


# ============================================================
# 第12页：给职校学生的建议
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), BLUE_DARK)
add_textbox(slide, Inches(0.5), Inches(0.25), Inches(12), Inches(0.8),
            "🌟 在AI时代，我们能做什么？", size=32, color=WHITE, bold=True)

tips = [
    ("📚", "学好专业技能", "你的专业 + AI = 更强的竞争力", "掌握本专业的核心技能，同时学会用AI工具提升效率"),
    ("🛠️", "学会使用AI工具", "让AI成为你的得力助手", "主动学习和体验各种AI工具，发现它们的用途"),
    ("💡", "培养AI做不到的能力", "创造力、沟通力、同理心", "这些是AI难以替代的人类独特能力"),
    ("🔍", "保持好奇心和终身学习", "技术一直在变，学习不能停", "关注AI发展，持续更新自己的知识和技能"),
]

for i, (icon, title, subtitle, desc) in enumerate(tips):
    y = Inches(1.5) + Inches(1.4) * i
    card = add_shape_bg(slide, Inches(1.0), y, Inches(11.3), Inches(1.2), RGBColor(0xF8, 0xF9, 0xFA))

    add_textbox(slide, Inches(1.3), y + Inches(0.1), Inches(0.8), Inches(0.8), icon, size=32, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(2.2), y + Inches(0.05), Inches(3), Inches(0.4), title, size=20, color=BLUE_DARK, bold=True)
    add_textbox(slide, Inches(2.2), y + Inches(0.5), Inches(3), Inches(0.5), subtitle, size=14, color=ACCENT_ORANGE, bold=True)
    add_textbox(slide, Inches(5.5), y + Inches(0.2), Inches(6.5), Inches(0.8), desc, size=15, color=MED_GRAY)

# 金句
add_shape_bg(slide, Inches(2), Inches(6.3), Inches(9.3), Inches(0.8), BLUE_DARK)
add_textbox(slide, Inches(2), Inches(6.35), Inches(9.3), Inches(0.7),
            "💬 不是AI取代你，而是先学会AI的人取代你。",
            size=20, color=WHITE, bold=True, align=PP_ALIGN.CENTER)


# ============================================================
# 第13页：课堂总结
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), BLUE_DARK)
add_textbox(slide, Inches(0.5), Inches(0.25), Inches(12), Inches(0.8),
            "📝 课堂总结", size=32, color=WHITE, bold=True)

summaries = [
    ("✅", "人工智能", "让机器模拟人类智能的技术\n感知→理解→推理→行动"),
    ("✅", "发展简史", "1956年诞生 → 多次起伏 → 当今快速发展"),
    ("✅", "广泛应用", "生活·学习·交通·医疗·娱乐·工业"),
    ("✅", "我们的态度", "拥抱AI · 学会协作 · 保持思考"),
]

for i, (icon, title, desc) in enumerate(summaries):
    x = Inches(0.8) + Inches(3.1) * i
    card = add_shape_bg(slide, x, Inches(1.6), Inches(2.8), Inches(3.5), RGBColor(0xE3, 0xF2, 0xFD))

    add_textbox(slide, x + Inches(0.1), Inches(1.8), Inches(2.6), Inches(0.5), icon + " " + title, size=20, color=BLUE_DARK, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, x + Inches(0.2), Inches(2.4), Inches(2.4), Inches(2.5), desc, size=15, color=DARK_GRAY, align=PP_ALIGN.CENTER)

# 思维导图式总结
add_shape_bg(slide, Inches(1), Inches(5.5), Inches(11.3), Inches(1.5), RGBColor(0x1A, 0x23, 0x7E))
add_textbox(slide, Inches(1), Inches(5.6), Inches(11.3), Inches(0.5),
            "核心思想", size=22, color=ACCENT_ORANGE, bold=True, align=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1), Inches(6.1), Inches(11.3), Inches(0.8),
            "人工智能不是遥远的未来，它已经来到我们身边。\n学习AI，不是为了成为AI专家，而是为了不被AI时代淘汰。",
            size=16, color=WHITE, align=PP_ALIGN.CENTER)


# ============================================================
# 第14页：课后任务
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), BLUE_DARK)
add_textbox(slide, Inches(0.5), Inches(0.25), Inches(12), Inches(0.8),
            "📋 课后任务", size=32, color=WHITE, bold=True)

# 必做
add_shape_bg(slide, Inches(1), Inches(1.6), Inches(5.2), Inches(4.5), RGBColor(0xE8, 0xF5, 0xE9))
add_textbox(slide, Inches(1.3), Inches(1.8), Inches(4.6), Inches(0.5), "✅ 必做任务", size=22, color=RGBColor(0x2E, 0x7D, 0x32), bold=True)

must_do = [
    "完成学习任务单",
    "探索一种新的AI工具",
    "（推荐：通义千问、文心一言、Kimi等）",
    "记录你的使用体验",
]
add_bullet_list(slide, Inches(1.3), Inches(2.4), Inches(4.6), Inches(3.0), must_do, size=16, color=DARK_GRAY)

# 选做
add_shape_bg(slide, Inches(7.0), Inches(1.6), Inches(5.3), Inches(4.5), RGBColor(0xFF, 0xF3, 0xE0))
add_textbox(slide, Inches(7.3), Inches(1.8), Inches(4.7), Inches(0.5), "⭐ 选做任务（任选其一）", size=20, color=ACCENT_ORANGE, bold=True)

nice_do = [
    "用AI帮你完成一次作业或学习任务",
    "采访一位家人对AI的看法",
    "了解一位AI领域的名人",
    "（图灵、李飞飞、吴恩达等）",
]
add_bullet_list(slide, Inches(7.3), Inches(2.4), Inches(4.7), Inches(3.0), nice_do, size=16, color=DARK_GRAY)


# ============================================================
# 第15页：结束页
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BLUE_DARK)

add_shape_bg(slide, Inches(0), Inches(2.5), Inches(13.333), Inches(0.04), ACCENT_ORANGE)

add_textbox(slide, Inches(1.5), Inches(3.0), Inches(10.3), Inches(1.5),
            "谢谢观看！", size=52, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

add_textbox(slide, Inches(2), Inches(4.5), Inches(9.3), Inches(0.8),
            "拥抱AI，拥抱未来 🚀", size=28, color=ACCENT_ORANGE, bold=True, align=PP_ALIGN.CENTER)

add_textbox(slide, Inches(2), Inches(5.5), Inches(9.3), Inches(0.6),
            "推荐AI工具：通义千问 · 文心一言 · Kimi · 智谱清言",
            size=16, color=RGBColor(0x90, 0xCA, 0xF9), align=PP_ALIGN.CENTER)

add_shape_bg(slide, Inches(0), Inches(6.5), Inches(13.333), Inches(0.04), ACCENT_ORANGE)


# ============================================================
# 保存
# ============================================================
output_path = "/home/admin/.openclaw/workspace/教学资料_初识人工智能/初识人工智能_教学PPT.pptx"
prs.save(output_path)
print(f"✅ PPT已生成：{output_path}")
