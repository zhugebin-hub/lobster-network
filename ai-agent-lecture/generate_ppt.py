#!/usr/bin/env python3
"""
生成「打造数字员工：AI 智能体全流程落地实战」PPT
授课时间：2026年6月6日 9:00-12:00
授课人：诸葛斌
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# 配色方案
DARK_BLUE = RGBColor(0x1B, 0x3A, 0x5C)      # 深蓝
MEDIUM_BLUE = RGBColor(0x2E, 0x75, 0xB6)     # 中蓝
LIGHT_BLUE = RGBColor(0x5B, 0x9B, 0xD5)      # 浅蓝
ACCENT_ORANGE = RGBColor(0xED, 0x7D, 0x31)   # 橙色强调
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
GREEN = RGBColor(0x70, 0xAD, 0x47)           # 绿色
RED = RGBColor(0xC0, 0x00, 0x00)             # 红色

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def add_background(slide, color):
    """设置幻灯片背景色"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color, border_color=None):
    """添加矩形形状"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.fill.solid()
        shape.line.fill.fore_color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height, text, font_size=18, 
                font_color=BLACK, bold=False, alignment=PP_ALIGN.LEFT,
                font_name='微软雅黑'):
    """添加文本框"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_rounded_rect(slide, left, top, width, height, fill_color):
    """添加圆角矩形"""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def create_title_slide():
    """创建封面页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
    add_background(slide, DARK_BLUE)
    
    # 顶部装饰条
    add_shape(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.15), ACCENT_ORANGE)
    
    # 主标题
    add_textbox(slide, Inches(1.5), Inches(2.0), Inches(10.3), Inches(1.5),
                '打造数字员工', font_size=44, font_color=WHITE, bold=True,
                alignment=PP_ALIGN.CENTER)
    
    add_textbox(slide, Inches(1.5), Inches(3.2), Inches(10.3), Inches(1.0),
                'AI 智能体全流程落地实战', font_size=36, font_color=LIGHT_BLUE, bold=True,
                alignment=PP_ALIGN.CENTER)
    
    # 分隔线
    add_shape(slide, Inches(4.5), Inches(4.5), Inches(4.3), Inches(0.04), ACCENT_ORANGE)
    
    # 主讲人信息
    add_textbox(slide, Inches(1.5), Inches(5.0), Inches(10.3), Inches(0.5),
                '诸葛斌', font_size=24, font_color=WHITE, bold=True,
                alignment=PP_ALIGN.CENTER)
    
    add_textbox(slide, Inches(1.5), Inches(5.5), Inches(10.3), Inches(0.5),
                '浙江工商大学 信息与电子工程学院 / 萨塞克斯人工智能学院',
                font_size=16, font_color=LIGHT_BLUE, alignment=PP_ALIGN.CENTER)
    
    add_textbox(slide, Inches(1.5), Inches(6.0), Inches(10.3), Inches(0.5),
                '2026年6月6日', font_size=16, font_color=LIGHT_BLUE,
                alignment=PP_ALIGN.CENTER)
    
    # 底部装饰条
    add_shape(slide, Inches(0), Inches(7.35), SLIDE_WIDTH, Inches(0.15), ACCENT_ORANGE)


def create_toc_slide():
    """创建目录页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, WHITE)
    
    # 顶部标题栏
    add_shape(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(1.0), DARK_BLUE)
    add_textbox(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.6),
                '课程目录', font_size=32, font_color=WHITE, bold=True)
    
    # 四个模块
    modules = [
        ('模块一', '为什么是现在？', '政策、趋势、必要性', '9:00-9:50'),
        ('模块二', '什么是数字员工？', '技术原理、架构、形态', '9:50-10:40'),
        ('模块三', '怎么落地？', '案例、路径、实操', '10:55-11:45'),
        ('模块四', '行动与 Q&A', '总结、行动清单、问答', '11:45-12:00'),
    ]
    
    for i, (num, title, desc, time) in enumerate(modules):
        top = Inches(1.5) + Inches(i * 1.4)
        
        # 模块编号
        shape = add_rounded_rect(slide, Inches(0.8), top, Inches(1.2), Inches(1.0), MEDIUM_BLUE)
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = num
        p.font.size = Pt(20)
        p.font.color.rgb = WHITE
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        
        # 标题
        add_textbox(slide, Inches(2.3), top + Inches(0.1), Inches(5), Inches(0.5),
                    title, font_size=24, font_color=DARK_BLUE, bold=True)
        
        # 描述
        add_textbox(slide, Inches(2.3), top + Inches(0.6), Inches(5), Inches(0.4),
                    desc, font_size=16, font_color=DARK_GRAY)
        
        # 时间
        add_textbox(slide, Inches(10.5), top + Inches(0.3), Inches(2), Inches(0.4),
                    time, font_size=18, font_color=ACCENT_ORANGE, bold=True,
                    alignment=PP_ALIGN.RIGHT)


def create_content_slide(title, content_blocks):
    """创建内容页（通用模板）"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, WHITE)
    
    # 顶部标题栏
    add_shape(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.9), DARK_BLUE)
    add_textbox(slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.6),
                title, font_size=28, font_color=WHITE, bold=True)
    
    # 内容区域
    current_top = Inches(1.2)
    for block in content_blocks:
        block_type = block.get('type', 'text')
        
        if block_type == 'text':
            add_textbox(slide, Inches(0.8), current_top, Inches(11.5), Inches(0.5),
                        block['text'], font_size=block.get('font_size', 18),
                        font_color=block.get('color', DARK_GRAY),
                        bold=block.get('bold', False))
            current_top += Inches(block.get('height', 0.5))
        
        elif block_type == 'bullet':
            for item in block['items']:
                add_textbox(slide, Inches(1.0), current_top, Inches(11), Inches(0.4),
                            f'• {item}', font_size=17, font_color=DARK_GRAY)
                current_top += Inches(0.45)
        
        elif block_type == 'quote':
            # 引用框
            quote_shape = add_rounded_rect(slide, Inches(0.8), current_top, Inches(11.5), 
                                           Inches(1.2), LIGHT_GRAY)
            add_textbox(slide, Inches(1.2), current_top + Inches(0.15), Inches(10.5), Inches(1.0),
                        block['text'], font_size=16, font_color=DARK_BLUE, bold=True)
            current_top += Inches(1.5)
        
        elif block_type == 'spacer':
            current_top += Inches(block.get('height', 0.3))
    
    return slide


def create_table_slide(title, headers, rows, col_widths=None):
    """创建表格页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, WHITE)
    
    # 顶部标题栏
    add_shape(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.9), DARK_BLUE)
    add_textbox(slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.6),
                title, font_size=28, font_color=WHITE, bold=True)
    
    # 计算表格尺寸
    num_cols = len(headers)
    num_rows = len(rows) + 1  # +1 for header
    table_width = Inches(11.5)
    row_height = Inches(0.5)
    table_height = Inches(num_rows * 0.45)
    
    if col_widths is None:
        col_width = table_width / num_cols
    else:
        total = sum(col_widths)
        col_width = table_width / total
    
    table = slide.shapes.add_table(num_rows, num_cols, Inches(0.8), Inches(1.3),
                                    table_width, table_height).table
    
    # 设置列宽
    if col_widths:
        total = sum(col_widths)
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(int(w / total * 11.5))
    else:
        col_w = int(table_width / num_cols)
        for i in range(num_cols):
            table.columns[i].width = col_w
    
    # 表头
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(14)
            paragraph.font.bold = True
            paragraph.font.color.rgb = WHITE
            paragraph.font.name = '微软雅黑'
            paragraph.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK_BLUE
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    # 数据行
    for r, row in enumerate(rows):
        for c, value in enumerate(row):
            cell = table.cell(r + 1, c)
            cell.text = str(value)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(12)
                paragraph.font.color.rgb = DARK_GRAY
                paragraph.font.name = '微软雅黑'
                paragraph.alignment = PP_ALIGN.CENTER
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT_GRAY if r % 2 == 0 else WHITE
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    return slide


def create_two_column_slide(title, left_title, left_content, right_title, right_content):
    """创建双栏对比页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, WHITE)
    
    # 顶部标题栏
    add_shape(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.9), DARK_BLUE)
    add_textbox(slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.6),
                title, font_size=28, font_color=WHITE, bold=True)
    
    # 左栏
    left_shape = add_rounded_rect(slide, Inches(0.5), Inches(1.3), Inches(5.8), Inches(5.5), LIGHT_GRAY)
    add_textbox(slide, Inches(0.8), Inches(1.5), Inches(5.2), Inches(0.5),
                left_title, font_size=22, font_color=DARK_BLUE, bold=True)
    
    for i, item in enumerate(left_content):
        add_textbox(slide, Inches(0.8), Inches(2.1) + Inches(i * 0.5), Inches(5.2), Inches(0.45),
                    f'• {item}', font_size=15, font_color=DARK_GRAY)
    
    # 右栏
    right_shape = add_rounded_rect(slide, Inches(7.0), Inches(1.3), Inches(5.8), Inches(5.5), LIGHT_GRAY)
    add_textbox(slide, Inches(7.3), Inches(1.5), Inches(5.2), Inches(0.5),
                right_title, font_size=22, font_color=DARK_BLUE, bold=True)
    
    for i, item in enumerate(right_content):
        add_textbox(slide, Inches(7.3), Inches(2.1) + Inches(i * 0.5), Inches(5.2), Inches(0.45),
                    f'• {item}', font_size=15, font_color=DARK_GRAY)


def create_architecture_slide():
    """创建架构图页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, WHITE)
    
    # 顶部标题栏
    add_shape(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.9), DARK_BLUE)
    add_textbox(slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.6),
                'AI Agent 标准架构', font_size=28, font_color=WHITE, bold=True)
    
    # 用户层
    user_shape = add_rounded_rect(slide, Inches(4.5), Inches(1.2), Inches(4.3), Inches(0.7), MEDIUM_BLUE)
    tf = user_shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = '👤 用 户'
    p.font.size = Pt(18)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    
    # 箭头
    add_textbox(slide, Inches(6.3), Inches(1.9), Inches(0.7), Inches(0.5),
                '▼', font_size=24, font_color=DARK_BLUE, bold=True, alignment=PP_ALIGN.CENTER)
    
    # 大脑层
    brain_shape = add_rounded_rect(slide, Inches(3.5), Inches(2.5), Inches(6.3), Inches(1.0), DARK_BLUE)
    tf = brain_shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = '🧠 大 脑：大 模 型（LLM）'
    p.font.size = Pt(20)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    
    # 说明文字
    add_textbox(slide, Inches(3.5), Inches(3.5), Inches(6.3), Inches(0.4),
                '理解意图 → 拆解任务 → 制定计划 → 决策下一步',
                font_size=14, font_color=DARK_GRAY, alignment=PP_ALIGN.CENTER)
    
    # 四个能力层
    capabilities = [
        ('👁️ 感知', '理解意图\n读取环境', LIGHT_BLUE),
        ('🧠 规划', '分解任务\n制定计划', MEDIUM_BLUE),
        ('🤲 行动', '调用工具\n执行操作', ACCENT_ORANGE),
        ('💾 记忆', '记住上下文\n积累经验', GREEN),
    ]
    
    for i, (name, desc, color) in enumerate(capabilities):
        left = Inches(0.8) + Inches(i * 3.1)
        cap_shape = add_rounded_rect(slide, left, Inches(4.3), Inches(2.8), Inches(1.5), color)
        tf = cap_shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = name
        p.font.size = Pt(16)
        p.font.color.rgb = WHITE
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(12)
        p2.font.color.rgb = WHITE
        p2.alignment = PP_ALIGN.CENTER
    
    # 底部说明
    add_textbox(slide, Inches(0.5), Inches(6.2), Inches(12), Inches(0.5),
                '这就是 OpenClaw（小龙虾）的架构 —— 也是数字员工的标准架构',
                font_size=16, font_color=DARK_BLUE, bold=True, alignment=PP_ALIGN.CENTER)


def create_flow_slide():
    """创建流程图页（落地四步走）"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, WHITE)
    
    # 顶部标题栏
    add_shape(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.9), DARK_BLUE)
    add_textbox(slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.6),
                '落地四步走战略', font_size=28, font_color=WHITE, bold=True)
    
    steps = [
        ('第一步', '试点\n1-3个月', '选择2-3个高价值场景\n内部验证', MEDIUM_BLUE),
        ('第二步', '扩展\n3-6个月', '扩展到10+场景\n建立数字员工矩阵', LIGHT_BLUE),
        ('第三步', '赋能\n6-12个月', '能力对外输出\n服务上下游企业', ACCENT_ORANGE),
        ('第四步', '生态\n12-24个月', '构建行业AI生态\n成为示范企业', GREEN),
    ]
    
    for i, (step, period, desc, color) in enumerate(steps):
        left = Inches(0.8) + Inches(i * 3.1)
        
        # 步骤框
        step_shape = add_rounded_rect(slide, left, Inches(1.5), Inches(2.8), Inches(2.5), color)
        tf = step_shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = step
        p.font.size = Pt(20)
        p.font.color.rgb = WHITE
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        p2 = tf.add_paragraph()
        p2.text = period
        p2.font.size = Pt(14)
        p2.font.color.rgb = WHITE
        p2.alignment = PP_ALIGN.CENTER
        
        # 描述
        add_textbox(slide, left, Inches(4.2), Inches(2.8), Inches(1.0),
                    desc, font_size=14, font_color=DARK_GRAY, alignment=PP_ALIGN.CENTER)
        
        # 箭头（除了最后一个）
        if i < 3:
            add_textbox(slide, left + Inches(2.8), Inches(2.3), Inches(0.3), Inches(0.5),
                        '→', font_size=28, font_color=DARK_BLUE, bold=True)
    
    # 底部总结
    add_shape(slide, Inches(0.5), Inches(5.8), Inches(12.3), Inches(1.2), LIGHT_GRAY)
    add_textbox(slide, Inches(0.8), Inches(5.9), Inches(11.5), Inches(1.0),
                '关键：先试点验证，再逐步扩展。不要一开始就追求完美，先跑通一个场景最重要。',
                font_size=16, font_color=DARK_BLUE, bold=True, alignment=PP_ALIGN.CENTER)


def create_case_slide(title, case_info):
    """创建案例页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, WHITE)
    
    # 顶部标题栏
    add_shape(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.9), DARK_BLUE)
    add_textbox(slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.6),
                title, font_size=28, font_color=WHITE, bold=True)
    
    y = Inches(1.3)
    for key, value in case_info.items():
        add_textbox(slide, Inches(0.8), y, Inches(2.5), Inches(0.4),
                    f'{key}：', font_size=16, font_color=DARK_BLUE, bold=True)
        add_textbox(slide, Inches(3.3), y, Inches(9), Inches(0.4),
                    value, font_size=16, font_color=DARK_GRAY)
        y += Inches(0.5)


def create_roi_slide():
    """创建 ROI 计算页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, WHITE)
    
    # 顶部标题栏
    add_shape(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.9), DARK_BLUE)
    add_textbox(slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.6),
                'ROI 计算 —— 数字员工的经济账', font_size=28, font_color=WHITE, bold=True)
    
    # 表格
    headers = ['项目', '传统模式', '智能体模式', '差异']
    rows = [
        ['人力成本', '1人×8万/年', '0.2人×8万/年+Token', '↓6.4万'],
        ['报告时间', '4小时/份', '15分钟/份', '↑16倍'],
        ['质量稳定性', '受状态影响', '始终如一', '↑'],
        ['Token成本', '0', '约5000元/年', '+0.5万'],
        ['年净节省', '—', '—', '≈6万/人'],
    ]
    
    table = slide.shapes.add_table(6, 4, Inches(1.5), Inches(1.3),
                                    Inches(10), Inches(2.8)).table
    
    table.columns[0].width = Inches(2.5)
    table.columns[1].width = Inches(2.5)
    table.columns[2].width = Inches(2.5)
    table.columns[3].width = Inches(2.5)
    
    # 表头
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK_BLUE
    
    # 数据
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r+1, c)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(13)
                p.font.color.rgb = DARK_GRAY
                p.alignment = PP_ALIGN.CENTER
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT_GRAY if r % 2 == 0 else WHITE
    
    # 结论框
    add_shape(slide, Inches(1.5), Inches(4.5), Inches(10), Inches(1.5), LIGHT_GRAY)
    add_textbox(slide, Inches(2.0), Inches(4.6), Inches(9), Inches(0.5),
                '💡 投资回收期：< 2个月', font_size=22, font_color=GREEN, bold=True,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(2.0), Inches(5.1), Inches(9), Inches(0.5),
                'Token 成本极低，但节省显著。数字员工的 ROI 远超传统 IT 投资。',
                font_size=16, font_color=DARK_BLUE, alignment=PP_ALIGN.CENTER)


def create_action_slide():
    """创建行动清单页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, WHITE)
    
    # 顶部标题栏
    add_shape(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.9), DARK_BLUE)
    add_textbox(slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.6),
                '行动清单', font_size=28, font_color=WHITE, bold=True)
    
    # 左栏：今天
    left_shape = add_rounded_rect(slide, Inches(0.5), Inches(1.3), Inches(5.8), Inches(5.5), LIGHT_GRAY)
    add_textbox(slide, Inches(0.8), Inches(1.5), Inches(5.2), Inches(0.5),
                '📌 今天回去就能做的3件事', font_size=20, font_color=DARK_BLUE, bold=True)
    
    today_items = [
        '识别场景——列出你工作中\n  最高频、最重复的3个任务',
        '体验 Agent——安装 OpenClaw\n  或试用 Coze/Dify，搭建\n  一个简单智能体',
        '组建团队——找1-2个同事，\n  组成 AI 探索小组',
    ]
    for i, item in enumerate(today_items):
        add_textbox(slide, Inches(0.8), Inches(2.2) + Inches(i * 1.5), Inches(5.2), Inches(1.3),
                    f'{i+1}. {item}', font_size=16, font_color=DARK_GRAY)
    
    # 右栏：本月
    right_shape = add_rounded_rect(slide, Inches(7.0), Inches(1.3), Inches(5.8), Inches(5.5), LIGHT_GRAY)
    add_textbox(slide, Inches(7.3), Inches(1.5), Inches(5.2), Inches(0.5),
                '📌 本月内能做的3件事', font_size=20, font_color=DARK_BLUE, bold=True)
    
    month_items = [
        '完成 MVP——选一个场景，\n  搭建数字员工原型',
        '收集数据——记录效率提升、\n  成本节省的具体数字',
        '向上汇报——用数据说服领导，\n  争取资源支持',
    ]
    for i, item in enumerate(month_items):
        add_textbox(slide, Inches(7.3), Inches(2.2) + Inches(i * 1.5), Inches(5.2), Inches(1.3),
                    f'{i+1}. {item}', font_size=16, font_color=DARK_GRAY)


def create_conclusion_slide():
    """创建核心结论页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, WHITE)
    
    # 顶部标题栏
    add_shape(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.9), DARK_BLUE)
    add_textbox(slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.6),
                '核心结论', font_size=28, font_color=WHITE, bold=True)
    
    conclusions = [
        ('国家战略', '国发〔2025〕11号：2027年智能体\n普及率≥70%，2030年≥90%\n这不是目标，是考核指标。', DARK_BLUE),
        ('技术成熟', 'Agent 架构已标准化：\n大脑+手脚+记忆+规划\nOpenClaw 等框架让搭建门槛大幅降低。', MEDIUM_BLUE),
        ('商业验证', '移动8万+电信50类，\nToken 计费模式跑通\n数字员工不是概念，是正在发生的现实。', ACCENT_ORANGE),
        ('行动窗口', '试错容错机制 + 国资考核变化\n+ 7万亿投资\n现在是入场的最佳时机。', GREEN),
    ]
    
    for i, (title, desc, color) in enumerate(conclusions):
        top = Inches(1.3) + Inches(i * 1.4)
        
        # 标题
        title_shape = add_rounded_rect(slide, Inches(0.8), top, Inches(2.5), Inches(0.6), color)
        tf = title_shape.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(18)
        p.font.color.rgb = WHITE
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        
        # 描述
        add_textbox(slide, Inches(3.5), top + Inches(0.05), Inches(8.5), Inches(0.6),
                    desc, font_size=15, font_color=DARK_GRAY)


def create_end_slide():
    """创建结束页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, DARK_BLUE)
    
    # 装饰条
    add_shape(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.15), ACCENT_ORANGE)
    
    # 感谢语
    add_textbox(slide, Inches(1.5), Inches(2.0), Inches(10.3), Inches(1.0),
                '感谢聆听！', font_size=48, font_color=WHITE, bold=True,
                alignment=PP_ALIGN.CENTER)
    
    # 分隔线
    add_shape(slide, Inches(4.5), Inches(3.3), Inches(4.3), Inches(0.04), ACCENT_ORANGE)
    
    # 金句
    add_shape(slide, Inches(2.0), Inches(3.8), Inches(9.3), Inches(1.5), RGBColor(0x15, 0x2D, 0x4A))
    add_textbox(slide, Inches(2.3), Inches(3.9), Inches(8.7), Inches(1.3),
                '"未来不属于会用 AI 的人，\n而属于会对话的人。"',
                font_size=24, font_color=LIGHT_BLUE, bold=True,
                alignment=PP_ALIGN.CENTER)
    
    # 联系方式
    add_textbox(slide, Inches(1.5), Inches(5.5), Inches(10.3), Inches(0.5),
                '诸葛斌 | 浙江工商大学', font_size=20, font_color=WHITE,
                alignment=PP_ALIGN.CENTER)
    
    add_textbox(slide, Inches(1.5), Inches(6.0), Inches(10.3), Inches(0.5),
                '课程群：打造数字员工AI智能体全流程落地实战', font_size=16, font_color=LIGHT_BLUE,
                alignment=PP_ALIGN.CENTER)
    
    # 底部装饰条
    add_shape(slide, Inches(0), Inches(7.35), SLIDE_WIDTH, Inches(0.15), ACCENT_ORANGE)


def main():
    """生成完整 PPT"""
    print('🦞 正在生成 PPT...')
    
    # P1 封面
    create_title_slide()
    
    # P2 目录
    create_toc_slide()
    
    # P3 开场问题
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, WHITE)
    add_shape(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.9), DARK_BLUE)
    add_textbox(slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.6),
                '今天，你用了几个 AI？', font_size=28, font_color=WHITE, bold=True)
    
    ai_items = ['早上查天气 → 用了 AI', '导航上班 → 用了 AI', '搜索信息 → 用了 AI', '翻译文档 → 用了 AI']
    for i, item in enumerate(ai_items):
        add_textbox(slide, Inches(1.5), Inches(1.5) + Inches(i * 0.6), Inches(10), Inches(0.5),
                    f'✅ {item}', font_size=22, font_color=DARK_GRAY)
    
    add_shape(slide, Inches(1.5), Inches(4.2), Inches(10.3), Inches(0.04), ACCENT_ORANGE)
    
    add_textbox(slide, Inches(1.5), Inches(4.5), Inches(10.3), Inches(0.6),
                '但这些都是"工具"，不是"员工"。', font_size=24, font_color=DARK_BLUE, bold=True,
                alignment=PP_ALIGN.CENTER)
    
    add_textbox(slide, Inches(1.5), Inches(5.3), Inches(10.3), Inches(0.4),
                '工具：你指挥它干活    vs    员工：你交代任务，它自己想办法完成',
                font_size=20, font_color=DARK_GRAY, alignment=PP_ALIGN.CENTER)
    
    add_textbox(slide, Inches(1.5), Inches(6.2), Inches(10.3), Inches(0.5),
                '今天的课程，就是从"用工具"到"管员工"的思维转变。',
                font_size=18, font_color=ACCENT_ORANGE, bold=True, alignment=PP_ALIGN.CENTER)
    
    # P4 三个信号
    create_table_slide('三个信号，同一个趋势',
                       ['信号', '来源', '含义'],
                       [
                           ['智能体普及率超70%', '国务院国发〔2025〕11号', '国家意志，2027年考核指标'],
                           ['移动上线8万数智员工', '2026数字中国峰会', '运营商实战，AI应用层争夺开始'],
                           ['电力系统投资4万亿', '2026年发改委', '能源底座，AI×能源双向赋能'],
                       ])
    
    # P5 三份政策文件
    create_table_slide('三份核心政策文件',
                       ['文件', '发文单位', '文号'],
                       [
                           ['国务院关于深入实施"人工智能+"行动的意见', '国务院', '国发〔2025〕11号'],
                           ['智能体规范应用与创新发展实施意见', '网信办、发改委、工信部', '配套文件'],
                           ['关于促进人工智能与能源双向赋能的行动方案', '发改委、能源局、工信部、数据局', '配套文件'],
                       ])
    
    # P6 70% 普及率原文
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, WHITE)
    add_shape(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.9), DARK_BLUE)
    add_textbox(slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.6),
                '70% 普及率 —— 国发〔2025〕11号原文', font_size=28, font_color=WHITE, bold=True)
    
    add_shape(slide, Inches(0.8), Inches(1.3), Inches(11.5), Inches(5.5), LIGHT_GRAY)
    
    quotes = [
        '到2027年，率先实现人工智能与6大重点领域广泛深度融合，',
        '新一代智能终端、智能体等应用普及率超70%',
        '',
        '到2030年，新一代智能终端、智能体等应用普及率超90%',
        '',
        '到2035年，我国全面步入智能经济和智能社会发展新阶段',
    ]
    for i, q in enumerate(quotes):
        if q:
            is_key = '70%' in q or '90%' in q
            add_textbox(slide, Inches(1.5), Inches(1.6) + Inches(i * 0.6), Inches(10), Inches(0.55),
                        q, font_size=20 if is_key else 18,
                        font_color=DARK_BLUE if is_key else DARK_GRAY,
                        bold=is_key, alignment=PP_ALIGN.CENTER)
    
    add_textbox(slide, Inches(0.8), Inches(6.5), Inches(11.5), Inches(0.5),
                '这是国务院正式文件，不是媒体报道，不是行业预测，是国家意志。',
                font_size=16, font_color=ACCENT_ORANGE, bold=True, alignment=PP_ALIGN.CENTER)
    
    # P7 四级考核体系
    create_table_slide('国家如何考核？——四级考核体系',
                       ['考核层级', '考核主体', '考核内容'],
                       [
                           ['一级：国务院统筹', '国家发展改革委', '统筹协调6大领域落实'],
                           ['二级：地方政府', '各省人民政府', '因地制宜贯彻落实'],
                           ['三级：国资考核', '国资委', '国有资本投资AI考核评价'],
                           ['四级：场景评价', '行业主管部门', '应用场景开放度评价'],
                       ])
    
    # P8 三个特殊设计
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, WHITE)
    add_shape(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.9), DARK_BLUE)
    add_textbox(slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.6),
                '考核机制的三个"特殊设计"', font_size=28, font_color=WHITE, bold=True)
    
    designs = [
        ('① 试错容错制度', '完善应用试错容错管理制度\n允许失败，不会因为试错而被问责。', LIGHT_BLUE),
        ('② 示范引领机制', '要强化示范引领，适时总结推广经验做法\n先试点、后推广。成功的案例全国推广。', ACCENT_ORANGE),
        ('③ 国资考核变化', '健全国有资本投资AI领域考核评价制度\n国企投资AI不再按传统ROI考核，\n而是按"战略价值+风险可控"考核。\n这意味着：国企可以"亏钱投AI"。', GREEN),
    ]
    
    for i, (title, desc, color) in enumerate(designs):
        top = Inches(1.3) + Inches(i * 1.9)
        title_shape = add_rounded_rect(slide, Inches(0.8), top, Inches(3.5), Inches(0.6), color)
        tf = title_shape.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(18)
        p.font.color.rgb = WHITE
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        
        add_textbox(slide, Inches(4.5), top + Inches(0.05), Inches(8), Inches(0.6),
                    desc, font_size=15, font_color=DARK_GRAY)
    
    # P9 6大领域
    create_table_slide('6大重点领域与智能体场景',
                       ['领域', '智能体典型场景'],
                       [
                           ['科学技术', '研发智能体、跨学科智能协作'],
                           ['产业发展', '工业全要素智能体、农业智能体、服务业智能体'],
                           ['消费提质', '智能助理、陪伴型智能体'],
                           ['民生福祉', '智能学伴、智能教师、健康助手、养老智能体'],
                           ['治理能力', '政务智能体、安全治理智能体'],
                           ['全球合作', '开源智能体、国际公共产品智能体'],
                       ])
    
    # P10 7万亿基建
    create_table_slide('7万亿基建投资 —— 硬件底座',
                       ['资金渠道', '规模', '用途'],
                       [
                           ['电力系统投资', '4万亿（2026年）', '绿电、新型电力系统'],
                           ['算力基础设施', '4000亿', '东数西算、智算中心'],
                           ['地下管网改造', '1万亿/年', 'AI传感器+实时监测'],
                           ['绿色债券+REITs', '数千亿', '绿电项目证券化'],
                           ['"两重两新"资金', '倾斜', '国家重大战略'],
                       ])
    
    add_textbox(prs.slides[-1], Inches(0.8), Inches(5.5), Inches(11.5), Inches(0.5),
                '没有这7万亿的硬件底座，70%普及率就是空话。',
                font_size=16, font_color=ACCENT_ORANGE, bold=True, alignment=PP_ALIGN.CENTER)
    
    # P11 AI × 能源 架构图
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, WHITE)
    add_shape(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.9), DARK_BLUE)
    add_textbox(slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.6),
                'AI × 能源 —— 不是加法，是乘法', font_size=28, font_color=WHITE, bold=True)
    
    # 智能体
    agent_shape = add_rounded_rect(slide, Inches(4.5), Inches(1.2), Inches(4.3), Inches(0.8), DARK_BLUE)
    tf = agent_shape.text_frame
    p = tf.paragraphs[0]
    p.text = '🦞 智能体（驾驶员）'
    p.font.size = Pt(18)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    
    add_textbox(slide, Inches(4.5), Inches(2.0), Inches(4.3), Inches(0.4),
                '70% 普及率 · 亿级并发 · 2027 目标', font_size=13, font_color=DARK_GRAY,
                alignment=PP_ALIGN.CENTER)
    
    # 双向箭头
    add_textbox(slide, Inches(5.5), Inches(2.5), Inches(2.3), Inches(0.5),
                '▼ 需要        反哺 ▲', font_size=16, font_color=ACCENT_ORANGE, bold=True,
                alignment=PP_ALIGN.CENTER)
    
    # 算力和能源
    compute_shape = add_rounded_rect(slide, Inches(1.5), Inches(3.2), Inches(4.5), Inches(1.5), MEDIUM_BLUE)
    tf = compute_shape.text_frame
    p = tf.paragraphs[0]
    p.text = '⚡ 算力（发动机）'
    p.font.size = Pt(18)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = '3700亿度/年  |  PUE<1.25'
    p2.font.size = Pt(14)
    p2.font.color.rgb = WHITE
    p2.alignment = PP_ALIGN.CENTER
    
    energy_shape = add_rounded_rect(slide, Inches(7.3), Inches(3.2), Inches(4.5), Inches(1.5), ACCENT_ORANGE)
    tf = energy_shape.text_frame
    p = tf.paragraphs[0]
    p.text = '🔋 能源（燃料）'
    p.font.size = Pt(18)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = '80% 绿电  |  新型电力系统'
    p2.font.size = Pt(14)
    p2.font.color.rgb = WHITE
    p2.alignment = PP_ALIGN.CENTER
    
    # 底部说明
    add_shape(slide, Inches(0.8), Inches(5.3), Inches(11.5), Inches(1.5), LIGHT_GRAY)
    add_textbox(slide, Inches(1.2), Inches(5.4), Inches(10.5), Inches(0.5),
                '能源 → AI：没有绿电，智能体跑不起来', font_size=18, font_color=DARK_BLUE, bold=True)
    add_textbox(slide, Inches(1.2), Inches(5.9), Inches(10.5), Inches(0.5),
                'AI → 能源：没有智能体，电网管不过来', font_size=18, font_color=DARK_BLUE, bold=True)
    add_textbox(slide, Inches(1.2), Inches(6.4), Inches(10.5), Inches(0.5),
                '这不是"AI+能源"，这是"AI×能源" —— 乘法关系，不是加法。',
                font_size=18, font_color=ACCENT_ORANGE, bold=True)
    
    # P12 模块一总结
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, DARK_BLUE)
    add_textbox(slide, Inches(1.5), Inches(2.0), Inches(10.3), Inches(0.8),
                '模块一小结', font_size=36, font_color=WHITE, bold=True,
                alignment=PP_ALIGN.CENTER)
    add_shape(slide, Inches(4.5), Inches(3.0), Inches(4.3), Inches(0.04), ACCENT_ORANGE)
    
    keywords = ['国家战略', '硬件底座', '行业机遇']
    descs = ['70%普及率是考核指标\n不是目标', '7万亿投资\n算力+绿电先行', '19个典型场景\n每个行业都有机会']
    
    for i, (kw, desc) in enumerate(zip(keywords, descs)):
        left = Inches(1.5) + Inches(i * 3.8)
        kw_shape = add_rounded_rect(slide, left, Inches(3.5), Inches(3.2), Inches(0.8), ACCENT_ORANGE)
        tf = kw_shape.text_frame
        p = tf.paragraphs[0]
        p.text = kw
        p.font.size = Pt(22)
        p.font.color.rgb = WHITE
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        
        add_textbox(slide, left, Inches(4.5), Inches(3.2), Inches(1.2),
                    desc, font_size=15, font_color=WHITE, alignment=PP_ALIGN.CENTER)
    
    add_textbox(slide, Inches(1.5), Inches(6.2), Inches(10.3), Inches(0.5),
                '问题留给模块三：你的行业，机会在哪里？',
                font_size=18, font_color=LIGHT_BLUE, alignment=PP_ALIGN.CENTER)
    
    # P13 三次进化
    create_table_slide('从 Chatbot 到 Agent —— 三次进化',
                       ['阶段', '代表', '能力', '局限'],
                       [
                           ['Chatbot\n聊天机器人', 'Siri、小爱同学', '问答、简单指令', '只能对话\n不能行动'],
                           ['Copilot\n副驾驶', 'GitHub Copilot\nOffice Copilot', '辅助创作\n代码补全', '需要人\n全程引导'],
                           ['Agent\n智能体/数字员工', 'OpenClaw\n各类AI Agent', '自主规划\n多步执行\n工具调用', '你交代任务\n它自己完成'],
                       ])
    
    # P14 Agent 定义
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, WHITE)
    add_shape(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.9), DARK_BLUE)
    add_textbox(slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.6),
                '什么是 AI Agent？', font_size=28, font_color=WHITE, bold=True)
    
    add_shape(slide, Inches(1.5), Inches(1.5), Inches(10.3), Inches(1.2), LIGHT_GRAY)
    add_textbox(slide, Inches(2.0), Inches(1.6), Inches(9.3), Inches(1.0),
                'AI Agent = 大模型的大脑 + 工具的手脚 + 记忆的经验 + 规划的思维',
                font_size=24, font_color=DARK_BLUE, bold=True, alignment=PP_ALIGN.CENTER)
    
    # P15 架构图
    create_architecture_slide()
    
    # P16 移动 vs 电信
    create_table_slide('运营商实战 —— 数字中国峰会最新数据',
                       ['维度', '中国移动', '中国电信'],
                       [
                           ['发布数量', '8万数智员工（已上线）', '8类数字员工（年底→50类）'],
                           ['应用方向', '向内（降本增效）', '向外（卖给客户）'],
                           ['应用场景', '网络鉴伪、研发设计\n营销服务、综合管理', '行政、财务、客服、销售'],
                           ['商业模式', '内部工具（节流）', '"数字员工+Token+连接"套餐'],
                           ['战略逻辑', '练内功 → 可能推向市场', '卖产品 → 直接创收'],
                       ])
    
    # P17 财务驱动
    create_table_slide('为什么两大运营商同时出手？',
                       ['运营商', '2026Q1营收', '同比', '2026Q1净利润', '同比'],
                       [
                           ['中国移动', '2665亿元', '+1.0%', '293亿元', '-4.2%'],
                           ['中国电信', '1314亿元', '+2.3%', '73.5亿元', '-17.1%'],
                       ])
    
    add_textbox(prs.slides[-1], Inches(0.8), Inches(5.0), Inches(11.5), Inches(0.5),
                '营收在涨，利润在跌。数字员工不是"锦上添花"，是"救命稻草"。',
                font_size=18, font_color=RED, bold=True, alignment=PP_ALIGN.CENTER)
    
    add_textbox(prs.slides[-1], Inches(0.8), Inches(5.6), Inches(11.5), Inches(0.5),
                '任何面临"成本端压力 > 收入端增长"的组织，数字员工都是必然选择。',
                font_size=16, font_color=DARK_BLUE, alignment=PP_ALIGN.CENTER)
    
    # P18 Token 计费
    create_table_slide('Token 计费 —— 智能体的"工资"',
                       ['传统计费', 'Token 计费', '差异'],
                       [
                           ['卖流量（按GB）', '卖Token（按消耗量）', 'Token是"智能"的计量单位'],
                           ['卖带宽（按Mbps）', '卖连接（按设备数）', '连接是"协同"的计量单位'],
                           ['卖软件（按License）', '卖数字员工（按任务数）', '数字员工是"劳动力"的计量单位'],
                       ])
    
    add_textbox(prs.slides[-1], Inches(0.8), Inches(4.8), Inches(11.5), Inches(0.5),
                'Token 消耗 = 智能体运行 = 持续收入',
                font_size=20, font_color=ACCENT_ORANGE, bold=True, alignment=PP_ALIGN.CENTER)
    
    # P19 三种形态
    create_table_slide('智能体的三种形态',
                       ['形态', '功能', '典型场景'],
                       [
                           ['调度智能体', '资源调度、成本优化', '电价低时多运行\n高时少运行'],
                           ['业务智能体', '业务流程自动化', '报告生成、数据整理\n客服应答'],
                           ['认知智能体', '复杂任务处理、决策', '综合调度决策\n安全风险评估'],
                       ])
    
    add_textbox(prs.slides[-1], Inches(0.8), Inches(4.5), Inches(11.5), Inches(0.5),
                '从简单到复杂，从自动化到智能化。',
                font_size=16, font_color=DARK_GRAY, alignment=PP_ALIGN.CENTER)
    
    # P20 虾尔能力展示
    create_table_slide('OpenClaw（小龙虾）—— 一个真实的数字员工',
                       ['能力类别', '具体能力', '实际案例'],
                       [
                           ['文档处理', '读/写/编辑/转换', '340+段落论文排版\nPPT分析评分'],
                           ['代码开发', 'Python/TS/Shell', 'AI知识问答系统\n课表转换脚本'],
                           ['定时任务', 'Cron 定时执行', '每日股市报告\n觅游社区互动'],
                           ['多平台通信', '钉钉/QQ/Telegram', '跨平台消息推送和响应'],
                           ['文件管理', '搜索/整理/打包', '12MB案例库自动整理'],
                           ['网络搜索', '联网搜索+信息整合', '政策研究、市场分析'],
                           ['协作网络', '多智能体协作', '龙虾网络\nNFS双向通道'],
                       ])
    
    add_textbox(prs.slides[-1], Inches(0.8), Inches(5.5), Inches(11.5), Inches(0.5),
                '这不是 PPT 上的概念，是每天在跑的真实工作。🦞',
                font_size=16, font_color=DARK_BLUE, bold=True, alignment=PP_ALIGN.CENTER)
    
    # P21 虾尔的一天
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, WHITE)
    add_shape(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.9), DARK_BLUE)
    add_textbox(slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.6),
                '虾尔的一天 —— 真实工作记录', font_size=28, font_color=WHITE, bold=True)
    
    timeline = [
        ('06:00', '🌅 起床（心跳轮询）', '检查龙虾网络消息、待处理任务\n觅游社区互动（点赞+评论）'),
        ('09:00', '📊 股市开盘', '获取持仓股票行情\n生成交易建议，推送钉钉群报告'),
        ('09:30', '💓 心跳检查', '邮件/日历/通知检查\n记忆文件整理'),
        ('10:00', '📝 课程建设', '101计划任务分解\n教学大纲编写、研讨会记录'),
        ('14:00', '🔧 项目开发', 'AI知识问答系统部署\n课表转换脚本开发、钉钉Bot对接'),
        ('19:00', '📚 文档处理', '论文排版、PPT分析\n案例整理'),
        ('21:30', '💓 晚间心跳', '社区互动、记忆归档'),
        ('00:00', '🌙 低功耗模式', '进入休眠，等待下次心跳'),
    ]
    
    for i, (time, activity, desc) in enumerate(timeline):
        top = Inches(1.2) + Inches(i * 0.7)
        
        # 时间
        add_textbox(slide, Inches(0.8), top, Inches(1.2), Inches(0.5),
                    time, font_size=16, font_color=ACCENT_ORANGE, bold=True)
        
        # 活动
        add_textbox(slide, Inches(2.0), top, Inches(3.5), Inches(0.5),
                    activity, font_size=15, font_color=DARK_BLUE, bold=True)
        
        # 描述
        add_textbox(slide, Inches(5.5), top, Inches(7), Inches(0.5),
                    desc, font_size=13, font_color=DARK_GRAY)
    
    add_textbox(slide, Inches(0.8), Inches(6.8), Inches(11.5), Inches(0.5),
                '7×24 小时，不请假、不抱怨、不跳槽。',
                font_size=16, font_color=DARK_BLUE, bold=True, alignment=PP_ALIGN.CENTER)
    
    # P22 数字员工 vs 传统软件
    create_table_slide('数字员工 vs 传统软件',
                       ['维度', '传统软件', '数字员工（Agent）'],
                       [
                           ['工作方式', '按固定流程执行', '理解意图，自主规划'],
                           ['适应能力', '需求变更=重新开发', '自然语言调整即可'],
                           ['学习成本', '需要培训操作', '会说话就会用'],
                           ['扩展性', '功能固定', '通过技能无限扩展'],
                           ['协作能力', '单机/局域网', '多智能体集群协作'],
                           ['24×7', '需要运维团队', '自主运行，自动恢复'],
                       ])
    
    add_textbox(prs.slides[-1], Inches(0.8), Inches(5.2), Inches(11.5), Inches(0.5),
                '数字员工不是"更好的软件"，是"另一种劳动力"。',
                font_size=18, font_color=ACCENT_ORANGE, bold=True, alignment=PP_ALIGN.CENTER)
    
    # P23 模块二总结
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, DARK_BLUE)
    add_textbox(slide, Inches(1.5), Inches(2.0), Inches(10.3), Inches(0.8),
                '模块二小结', font_size=36, font_color=WHITE, bold=True,
                alignment=PP_ALIGN.CENTER)
    add_shape(slide, Inches(4.5), Inches(3.0), Inches(4.3), Inches(0.04), ACCENT_ORANGE)
    
    keywords2 = ['Agent ≠ Chatbot', '标准架构', '新劳动力']
    descs2 = ['能自主规划\n执行、记忆', '大脑+手脚\n+记忆+规划', '不是软件升级\n是劳动力形态变革']
    
    for i, (kw, desc) in enumerate(zip(keywords2, descs2)):
        left = Inches(1.5) + Inches(i * 3.8)
        kw_shape = add_rounded_rect(slide, left, Inches(3.5), Inches(3.2), Inches(0.8), ACCENT_ORANGE)
        tf = kw_shape.text_frame
        p = tf.paragraphs[0]
        p.text = kw
        p.font.size = Pt(20)
        p.font.color.rgb = WHITE
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        
        add_textbox(slide, left, Inches(4.5), Inches(3.2), Inches(1.2),
                    desc, font_size=15, font_color=WHITE, alignment=PP_ALIGN.CENTER)
    
    add_textbox(slide, Inches(1.5), Inches(6.2), Inches(10.3), Inches(0.5),
                '下一模块：怎么把数字员工落到你的组织里？',
                font_size=18, font_color=LIGHT_BLUE, alignment=PP_ALIGN.CENTER)
    
    # P24 落地四步走
    create_flow_slide()
    
    # P25 试点场景选择
    create_table_slide('第一步 —— 试点场景选择',
                       ['标准', '说明', '举例'],
                       [
                           ['高价值', '能显著降本增效', '巡检、预警、报告'],
                           ['低风险', '失败影响可控', '文档整理、数据汇总'],
                           ['易验证', '3个月内可出效果', '报告生成效率提升'],
                           ['可复制', '能推广到其他场景', '巡检→维护→安全'],
                       ])
    
    add_textbox(prs.slides[-1], Inches(0.8), Inches(4.5), Inches(11.5), Inches(0.5),
                '建议：选一个"痛点明显、数据可得、流程清晰"的场景开始。',
                font_size=16, font_color=DARK_BLUE, bold=True, alignment=PP_ALIGN.CENTER)
    
    # P26 案例1：教育
    create_case_slide('实战案例 1 —— 教育行业', {
        '场景': '信电学院课程知识问答',
        '痛点': '学生问题重复、教师答疑耗时',
        '方案': 'OpenClaw + 本地RAG + 钉钉Bot',
        '知识库': '6门课程、77个知识块',
        '效果': '7×24小时自动答疑，教师答疑工作量↓70%',
        '架构': '知识库层 → 智能体层 → 应用交互层',
    })
    
    # P27 案例2：课程评价
    create_table_slide('实战案例 2 —— 课程期末评价',
                       ['维度', '传统方式', 'AI辅助方式', '效率提升'],
                       [
                           ['PPT阅读', '教师逐份阅读\n（约20小时）', 'AI批量提取+分析\n（约30分钟）', '40倍'],
                           ['评分', '主观打分\n标准不一', '五维评分标准\n一致性强', '—'],
                           ['报告生成', '手动编写\n（约8小时）', 'AI生成初稿+人工调整\n（约1小时）', '8倍'],
                           ['总耗时', '约30小时', '约2小时', '15倍'],
                       ])
    
    add_textbox(prs.slides[-1], Inches(0.8), Inches(5.0), Inches(11.5), Inches(0.5),
                '产出：评分表 + 教学分析报告 + 汇报PPT + 教学案例包',
                font_size=16, font_color=DARK_BLUE, bold=True, alignment=PP_ALIGN.CENTER)
    
    # P28 案例3：课程建设
    create_case_slide('实战案例 3 —— 101计划课程建设', {
        '项目': '《计算机网络》101计划课程建设',
        '团队': '5人（教授+副教授+实验师）',
        'AI角色': '任务分解、进度跟踪、文档整理、会议纪要',
        '产出': '任务分工表 + 研讨会记录 + 6月行动计划',
        '效果': '原本需要一周的规划工作，一天完成',
        '启示': 'AI不是替代团队，是让团队聚焦于"只有人能做的事"',
    })
    
    # P29 多行业应用
    create_table_slide('数字员工的多行业应用',
                       ['行业', '场景', '数字员工类型', '效果'],
                       [
                           ['教育', '作业批改、课程PPT、答疑', '业务智能体', '效率↑10倍'],
                           ['能源', '管道巡检、故障预警\n泵站调度', '业务+调度智能体', '成本↓30%'],
                           ['金融', '报告生成、数据分析、风控', '认知智能体', '耗时↓80%'],
                           ['政务', '政策咨询、材料审核\n数据汇总', '业务智能体', '响应↓70%'],
                           ['制造', '设备维护、质量检测\n排产调度', '调度+业务智能体', '故障↓50%'],
                       ])
    
    add_textbox(prs.slides[-1], Inches(0.8), Inches(5.0), Inches(11.5), Inches(0.5),
                '每个行业都有高频、重复、规则清晰的场景 —— 这些就是数字员工的用武之地。',
                font_size=16, font_color=DARK_BLUE, bold=True, alignment=PP_ALIGN.CENTER)
    
    # P30 ROI 计算
    create_roi_slide()
    
    # P31 技术选型
    create_table_slide('技术选型 —— 用什么搭建数字员工？',
                       ['方案', '适合场景', '门槛', '灵活性'],
                       [
                           ['OpenClaw\n（小龙虾）', '通用、多平台\n多技能', '中', '⭐⭐⭐⭐⭐'],
                           ['阿里百炼平台', '企业级\n阿里云生态', '中', '⭐⭐⭐⭐'],
                           ['Dify', '快速原型\n可视化编排', '低', '⭐⭐⭐'],
                           ['LangChain', '开发者\n深度定制', '高', '⭐⭐⭐⭐⭐'],
                           ['Coze（扣子）', '个人/小团队\n快速上手', '低', '⭐⭐'],
                       ])
    
    add_textbox(prs.slides[-1], Inches(0.8), Inches(5.0), Inches(11.5), Inches(0.5),
                '推荐：OpenClaw —— 技能生态丰富、多平台通信、支持多智能体协作。',
                font_size=16, font_color=DARK_BLUE, bold=True, alignment=PP_ALIGN.CENTER)
    
    # P32 行动清单
    create_action_slide()
    
    # P33 核心结论
    create_conclusion_slide()
    
    # P34 Q&A
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, WHITE)
    add_shape(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.9), DARK_BLUE)
    add_textbox(slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.6),
                'Q & A', font_size=28, font_color=WHITE, bold=True)
    
    add_textbox(slide, Inches(1.5), Inches(3.0), Inches(10.3), Inches(1.0),
                '欢迎提问！', font_size=36, font_color=DARK_BLUE, bold=True,
                alignment=PP_ALIGN.CENTER)
    
    # P35 结束页
    create_end_slide()
    
    # 保存
    output_path = '/home/admin/.openclaw/workspace/ai-agent-lecture/打造数字员工-AI智能体全流程落地实战.pptx'
    prs.save(output_path)
    print(f'✅ PPT 已保存：{output_path}')
    print(f'📊 总页数：{len(prs.slides)}')


if __name__ == '__main__':
    main()
