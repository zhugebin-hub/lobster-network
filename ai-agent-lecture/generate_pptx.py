#!/usr/bin/env python3
"""
生成浙大授课定制版PPTX文件
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

def create_presentation():
    """创建PPT演示文稿"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # 定义颜色方案
    colors = {
        'primary': RGBColor(0x1E, 0x88, 0xE5),  # 蓝色
        'secondary': RGBColor(0x43, 0xA0, 0x47),  # 绿色
        'accent': RGBColor(0xFF, 0x6F, 0x00),  # 橙色
        'dark': RGBColor(0x26, 0x32, 0x38),  # 深灰
        'light': RGBColor(0xF5, 0xF5, 0xF5),  # 浅灰
        'white': RGBColor(0xFF, 0xFF, 0xFF),
        'text': RGBColor(0x21, 0x21, 0x21),
        'text_light': RGBColor(0x75, 0x75, 0x75),
    }
    
    return prs, colors

def add_title_slide(prs, colors):
    """添加封面幻灯片"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
    
    # 背景
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = colors['primary']
    
    # 主标题
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11.333), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "打造数字员工：AI智能体（Agent）全流程落地实战"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = colors['white']
    p.alignment = PP_ALIGN.CENTER
    
    # 副标题
    txBox = slide.shapes.add_textbox(Inches(1), Inches(3.2), Inches(11.333), Inches(1))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "浙江大学—国家管网集团甘肃分公司人工智能实战培训班"
    p.font.size = Pt(24)
    p.font.color.rgb = colors['white']
    p.alignment = PP_ALIGN.CENTER
    
    # 讲师信息
    txBox = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(11.333), Inches(0.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "主讲人：诸葛斌"
    p.font.size = Pt(20)
    p.font.color.rgb = colors['white']
    p.alignment = PP_ALIGN.CENTER
    
    # 日期和地点
    txBox = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(11.333), Inches(0.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "2026年5月28日 | 浙江大学华家池校区"
    p.font.size = Pt(18)
    p.font.color.rgb = colors['white']
    p.alignment = PP_ALIGN.CENTER
    
    # 底部装饰线
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 
        Inches(4), Inches(6.5), Inches(5.333), Inches(0.05)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = colors['white']
    shape.line.fill.background()

def add_content_slide(prs, colors, title, content_blocks):
    """添加内容幻灯片"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 背景
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = colors['light']
    
    # 标题栏
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(13.333), Inches(1)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = colors['primary']
    shape.line.fill.background()
    
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(12.333), Inches(0.7))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = colors['white']
    
    # 内容区域
    y_offset = Inches(1.3)
    
    for block in content_blocks:
        block_type = block.get('type', 'text')
        
        if block_type == 'text':
            txBox = slide.shapes.add_textbox(
                Inches(0.8), y_offset, Inches(11.733), Inches(0.6)
            )
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = block['text']
            p.font.size = Pt(18)
            p.font.color.rgb = colors['text']
            y_offset += Inches(0.7)
            
        elif block_type == 'bullet':
            txBox = slide.shapes.add_textbox(
                Inches(1), y_offset, Inches(11.333), Inches(0.5)
            )
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = f"• {block['text']}"
            p.font.size = Pt(16)
            p.font.color.rgb = colors['text']
            y_offset += Inches(0.55)
            
        elif block_type == 'highlight':
            # 高亮框
            shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(0.8), y_offset, Inches(11.733), Inches(0.8)
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(0xE3, 0xF2, 0xFD)
            shape.line.color.rgb = colors['primary']
            shape.line.width = Pt(1)
            
            txBox = slide.shapes.add_textbox(
                Inches(1.2), y_offset + Inches(0.1), Inches(11.333), Inches(0.6)
            )
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = block['text']
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = colors['primary']
            y_offset += Inches(0.9)
            
        elif block_type == 'quote':
            txBox = slide.shapes.add_textbox(
                Inches(1.5), y_offset, Inches(10.333), Inches(0.8)
            )
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = f"> {block['text']}"
            p.font.size = Pt(16)
            p.font.italic = True
            p.font.color.rgb = colors['text_light']
            y_offset += Inches(0.9)
            
        elif block_type == 'table':
            rows = block.get('rows', [])
            cols = block.get('cols', 4)
            if rows:
                table_shape = slide.shapes.add_table(
                    len(rows), cols,
                    Inches(0.8), y_offset, Inches(11.733), Inches(0.4 * len(rows))
                )
                table = table_shape.table
                
                for i, row_data in enumerate(rows):
                    for j, cell_text in enumerate(row_data):
                        cell = table.cell(i, j)
                        cell.text = str(cell_text)
                        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                        
                        for paragraph in cell.text_frame.paragraphs:
                            paragraph.font.size = Pt(12)
                            if i == 0:
                                paragraph.font.bold = True
                                paragraph.font.color.rgb = colors['white']
                            else:
                                paragraph.font.color.rgb = colors['text']
                
                y_offset += Inches(0.4 * len(rows) + 0.2)
    
    return slide

def add_conclusion_slide(prs, colors):
    """添加结论幻灯片"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = colors['primary']
    
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(11.333), Inches(1))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "核心结论"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = colors['white']
    p.alignment = PP_ALIGN.CENTER
    
    conclusions = [
        "国发〔2025〕11号文件明确：2027年智能体普及率≥70%",
        "7万亿基建投资是硬件底座，19个典型场景是软件入口",
        "两大运营商同时出手，AI应用层争夺已开始",
        "管网集团机会：对内降本增效，对外输出创收",
        "率先打通'算力-绿电-智能体-Token'链路的企业将占先机"
    ]
    
    y = Inches(2.5)
    for i, text in enumerate(conclusions):
        txBox = slide.shapes.add_textbox(Inches(1.5), y, Inches(10.333), Inches(0.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"{i+1}. {text}"
        p.font.size = Pt(20)
        p.font.color.rgb = colors['white']
        y += Inches(0.6)
    
    # 底部
    txBox = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(11.333), Inches(0.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "🦞 小龙虾就是这条流水线上最先跑起来的那只龙虾"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = colors['white']
    p.alignment = PP_ALIGN.CENTER

def add_qa_slide(prs, colors):
    """添加Q&A幻灯片"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = colors['dark']
    
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.333), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Q&A"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = colors['white']
    p.alignment = PP_ALIGN.CENTER
    
    txBox = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(11.333), Inches(1))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "谢谢！"
    p.font.size = Pt(36)
    p.font.color.rgb = colors['white']
    p.alignment = PP_ALIGN.CENTER

def main():
    prs, colors = create_presentation()
    
    # 1. 封面
    add_title_slide(prs, colors)
    
    # 2. 课程导引
    add_content_slide(prs, colors, "课程导引——为什么是现在？", [
        {'type': 'highlight', 'text': '三个信号，同一个趋势'},
        {'type': 'table', 'rows': [
            ['信号', '来源', '含义'],
            ['"智能体普及率超70%"', '国务院国发〔2025〕11号', '国家意志，2027年考核指标'],
            ['移动上线8万数智员工', '2026数字中国峰会', '运营商实战，AI应用层争夺开始'],
            ['电力系统投资4万亿', '2026年发改委', '能源底座，AI×能源双向赋能'],
        ], 'cols': 3},
        {'type': 'quote', 'text': '这不是技术趋势，是国家战略。不是"要不要做"，是"多快做好"。'},
    ])
    
    # 3. 课程目标
    add_content_slide(prs, colors, "课程目标", [
        {'type': 'highlight', 'text': '学完本课程，你将能够：'},
        {'type': 'bullet', 'text': '理解政策——国发〔2025〕11号对能源行业的核心要求'},
        {'type': 'bullet', 'text': '掌握技术——AI智能体的三种形态与落地路径'},
        {'type': 'bullet', 'text': '对标实战——移动8万数智员工的经验教训'},
        {'type': 'bullet', 'text': '设计场景——为国家管网集团定制数字员工方案'},
        {'type': 'bullet', 'text': '计算ROI——Token计费模式下的成本收益分析'},
    ])
    
    # 4. 政策面
    add_content_slide(prs, colors, "第一部分：政策面——国发〔2025〕11号深度解读", [
        {'type': 'highlight', 'text': '三份核心政策文件'},
        {'type': 'table', 'rows': [
            ['文件', '发文单位', '文号'],
            ['《国务院关于深入实施"人工智能+"行动的意见》', '国务院', '国发〔2025〕11号'],
            ['《智能体规范应用与创新发展实施意见》', '网信办、发改委、工信部', '配套文件'],
            ['《关于促进人工智能与能源双向赋能的行动方案》', '发改委、能源局、工信部、数据局', '配套文件'],
        ], 'cols': 3},
        {'type': 'quote', 'text': '国发11号是顶层设计，另外两份是专项落地文件。'},
    ])
    
    # 5. 70%普及率
    add_content_slide(prs, colors, "70%普及率的官方表述", [
        {'type': 'quote', 'text': '"到2027年，率先实现人工智能与6大重点领域广泛深度融合，新一代智能终端、智能体等应用普及率超70%"'},
        {'type': 'quote', 'text': '"到2030年，新一代智能终端、智能体等应用普及率超90%"'},
        {'type': 'quote', 'text': '"到2035年，我国全面步入智能经济和智能社会发展新阶段"'},
        {'type': 'highlight', 'text': '这是国务院正式文件，不是媒体报道，不是行业预测，是国家意志。'},
    ])
    
    # 6. 6大重点领域
    add_content_slide(prs, colors, "6大重点领域与能源行业的关联", [
        {'type': 'table', 'rows': [
            ['领域', '智能体应用场景', '管网集团关联'],
            ['产业发展', '工业全要素智能体', '⭐⭐⭐⭐⭐ 管道巡检、调度优化、设备维护'],
            ['治理能力', '安全治理智能体', '⭐⭐⭐⭐ 安全生产监管、防灾减灾'],
            ['民生福祉', '健康助手、养老智能体', '⭐⭐ 员工健康管理'],
            ['科学技术', '研发智能体', '⭐⭐⭐ 管网技术攻关'],
            ['消费提质', '智能助理', '⭐ 办公自动化'],
            ['全球合作', '开源智能体', '⭐ 技术生态建设'],
        ], 'cols': 3},
    ])
    
    # 7. 考核体系
    add_content_slide(prs, colors, "国家如何考核？——四级考核体系", [
        {'type': 'table', 'rows': [
            ['考核层级', '考核主体', '考核内容', '管网集团应对'],
            ['一级：国务院统筹', '国家发展改革委', '统筹协调6大领域落实', '对接国家AI+能源政策'],
            ['二级：地方政府', '各省人民政府', '因地制宜贯彻落实', '配合甘肃省考核要求'],
            ['三级：国资考核', '国资委', '国有资本投资AI考核评价', '战略价值>ROI，可"亏钱投AI"'],
            ['四级：场景评价', '行业主管部门', '应用场景开放度评价', '开放管网场景，争取示范'],
        ], 'cols': 4},
    ])
    
    # 8. 考核特殊设计
    add_content_slide(prs, colors, "考核机制的特殊设计", [
        {'type': 'highlight', 'text': '① "试错容错"制度'},
        {'type': 'quote', 'text': '"完善应用试错容错管理制度"——允许失败，不会因为试错而被问责。'},
        {'type': 'highlight', 'text': '② "示范引领"机制'},
        {'type': 'quote', 'text': '"要强化示范引领，适时总结推广经验做法"——先试点、后推广。管网集团可争取成为"示范企业"。'},
        {'type': 'highlight', 'text': '③ "国有资本考核"变化'},
        {'type': 'quote', 'text': '"健全国有资本投资人工智能领域考核评价和风险监管等制度"——战略价值>ROI。'},
    ])
    
    # 9. 能源面
    add_content_slide(prs, colors, "第二部分：能源面——AI+能源双向赋能", [
        {'type': 'highlight', 'text': '核心逻辑链'},
        {'type': 'text', 'text': '能源 → AI：没有绿电，智能体跑不起来（3700亿度/年）'},
        {'type': 'text', 'text': 'AI → 能源：没有智能体，电网管不过来（六网协同需要AI调度）'},
        {'type': 'highlight', 'text': '这不是"AI+能源"，这是"AI×能源"——乘法关系，不是加法。'},
    ])
    
    # 10. 用电指标
    add_content_slide(prs, colors, "用电硬性指标", [
        {'type': 'table', 'rows': [
            ['指标', '数据', '对管网集团的含义'],
            ['2025年数据中心用电', '3700亿千瓦时', '算力需求激增'],
            ['新建算力绿电要求', '≥80%', '管网数据中心需绿电改造'],
            ['PUE上限', '1.25以下', '机房能效必须达标'],
            ['AI用电占全国比例', '3-5%（未来）', '能源企业是AI用电大户'],
        ], 'cols': 3},
        {'type': 'highlight', 'text': '管网集团的泵站、压缩机站、数据中心，都是"算力+能源"的交汇点。'},
    ])
    
    # 11. AI反哺能源
    add_content_slide(prs, colors, "AI反哺能源的四个层次", [
        {'type': 'table', 'rows': [
            ['层次', '管网集团应用场景', '效率提升', '智能体角色'],
            ['① 算力帮能源省钱', '电价低时多运行、高时少运行', '用电成本↓7-15%', '调度智能体'],
            ['② 智能体改造业务', '管道巡检、故障预警', '故障处置↓80%', '业务智能体'],
            ['③ 能源交易智能体', '电力采购优化、绿电交易', '综合成本↓7-15%', '交易智能体'],
            ['④ 大模型+智能体协同', '综合调度平台', '任务耗时↓80%', '认知智能体'],
        ], 'cols': 4},
        {'type': 'quote', 'text': '雄安电网调度智能体：故障处置时长压缩80%'},
        {'type': 'quote', 'text': '上海嘉定监控智能体：报告生成效率提升15倍'},
    ])
    
    # 12. 资金通道
    add_content_slide(prs, colors, "资金通道：7万亿与管网集团", [
        {'type': 'table', 'rows': [
            ['资金渠道', '政策工具', '管网集团可申请'],
            ['电力系统投资', '2026年4万亿', '泵站/压缩机站绿电改造'],
            ['算力基础设施', '东数西算+智算中心', '数据中心升级'],
            ['基础设施REITs', '绿电项目证券化', '新能源项目融资'],
            ['绿色债券', '碳中和债券', '零碳场站建设'],
            ['"两重两新"资金', '国家重大战略', 'AI+管网示范工程'],
        ], 'cols': 3},
    ])
    
    # 13. 实战面
    add_content_slide(prs, colors, "第三部分：实战面——运营商双雄对决", [
        {'type': 'highlight', 'text': '数字中国峰会最新数据'},
        {'type': 'table', 'rows': [
            ['维度', '中国移动', '中国电信'],
            ['发布数量', '8万数智员工（已上线）', '8类数字员工（年底→50类）'],
            ['应用方向', '向内（降本增效）', '向外（卖给客户）'],
            ['应用场景', '网络鉴伪、研发设计、营销服务', '行政、财务、客服、销售'],
            ['商业模式', '内部工具（节流）', '"数字员工+Token+连接"套餐'],
        ], 'cols': 3},
    ])
    
    # 14. 财务压力
    add_content_slide(prs, colors, "运营商财务压力驱动", [
        {'type': 'table', 'rows': [
            ['运营商', '2026Q1营收', '同比', '2026Q1净利润', '同比'],
            ['中国移动', '2665亿元', '+1.0%', '293亿元', '-4.2%'],
            ['中国电信', '1314亿元', '+2.3%', '73.5亿元', '-17.1%'],
        ], 'cols': 5},
        {'type': 'highlight', 'text': '营收在涨，利润在跌。数字员工不是"锦上添花"，是"救命稻草"。'},
        {'type': 'quote', 'text': '管网集团同样面临：成本端压力 > 收入端增长。数字员工是必然选择。'},
    ])
    
    # 15. 两种打法
    add_content_slide(prs, colors, "两种打法，两种逻辑", [
        {'type': 'highlight', 'text': '中国移动：向内降本（管网集团可借鉴）'},
        {'type': 'bullet', 'text': '8万数智员工全部用在自己身上'},
        {'type': 'bullet', 'text': '网络鉴伪、研发设计、营销服务、综合管理'},
        {'type': 'bullet', 'text': '启示：管网集团可先内部试点，再推向市场'},
        {'type': 'highlight', 'text': '中国电信：向外创收（管网集团可借鉴）'},
        {'type': 'bullet', 'text': '面向中小微企业发布数字员工套餐'},
        {'type': 'bullet', 'text': '"数字员工+Token+连接"，低成本、快速应用'},
        {'type': 'bullet', 'text': '启示：管网集团可将AI能力打包，服务上下游企业'},
    ])
    
    # 16. Token计费
    add_content_slide(prs, colors, "Token计费的战略意义", [
        {'type': 'quote', 'text': '国发〔2025〕11号："发展模型即服务、智能体即服务等，打造人工智能应用服务链。"'},
        {'type': 'highlight', 'text': '电信的"数字员工+Token+连接"套餐，就是"智能体即服务"的实战落地。'},
        {'type': 'table', 'rows': [
            ['传统计费', 'Token计费', '差异'],
            ['卖流量（按GB）', '卖Token（按消耗量）', 'Token是"智能"的计量单位'],
            ['卖带宽（按Mbps）', '卖连接（按设备数）', '连接是"协同"的计量单位'],
            ['卖软件（按License）', '卖数字员工（按任务数）', '数字员工是"劳动力"的计量单位'],
        ], 'cols': 3},
        {'type': 'highlight', 'text': 'Token消耗 = 智能体运行 = 持续收入。'},
    ])
    
    # 17. 技术面
    add_content_slide(prs, colors, "第四部分：技术面——智能体=数字员工", [
        {'type': 'highlight', 'text': '智能体的三种形态'},
        {'type': 'table', 'rows': [
            ['形态', '功能', '管网集团案例'],
            ['调度智能体', '资源调度、成本优化', '泵站运行调度、电力采购优化'],
            ['业务智能体', '业务流程自动化', '管道巡检、故障预警、报告生成'],
            ['认知智能体', '复杂任务处理、决策', '综合调度决策、安全风险评估'],
        ], 'cols': 3},
    ])
    
    # 18. 映射关系
    add_content_slide(prs, colors, "智能体 = 数字员工：完整的映射", [
        {'type': 'table', 'rows': [
            ['维度', '政策语言', '运营商语言', '管网集团语言'],
            ['形态', '智能体（AI Agent）', '数字员工/数智员工', '管网数字员工'],
            ['数量', '70%普及率 = 亿级', '移动8万+电信50类', '管网数字员工矩阵'],
            ['算力', '3700亿度/年', 'Token计费', '泵站/数据中心算力'],
            ['调度', '算力帮电力调峰', '网络鉴伪/客服智能体', '管道调度/巡检智能体'],
            ['应用', '19个典型场景', '8万内部场景+50类外部场景', '管网10大核心场景'],
        ], 'cols': 4},
    ])
    
    # 19. 10大场景
    add_content_slide(prs, colors, "管网集团10大核心场景设计", [
        {'type': 'table', 'rows': [
            ['场景', '智能体类型', '预期效果', '优先级'],
            ['① 管道巡检', '业务智能体', '巡检效率↑300%，漏检率↓90%', '⭐⭐⭐⭐⭐'],
            ['② 故障预警', '认知智能体', '故障提前72小时预警', '⭐⭐⭐⭐⭐'],
            ['③ 泵站调度', '调度智能体', '用电成本↓15%', '⭐⭐⭐⭐'],
            ['④ 设备维护', '业务智能体', '维护成本↓30%', '⭐⭐⭐⭐'],
            ['⑤ 安全监控', '业务智能体', '安全事故↓80%', '⭐⭐⭐⭐⭐'],
            ['⑥ 报告生成', '业务智能体', '报告时间↓90%', '⭐⭐⭐'],
            ['⑦ 电力交易', '交易智能体', '用电成本↓10%', '⭐⭐⭐'],
            ['⑧ 客户服务', '业务智能体', '响应时间↓70%', '⭐⭐⭐'],
            ['⑨ 培训教育', '认知智能体', '培训效率↑200%', '⭐⭐'],
            ['⑩ 综合决策', '认知智能体', '决策效率↑50%', '⭐⭐⭐⭐'],
        ], 'cols': 4},
    ])
    
    # 20. 落地路径
    add_content_slide(prs, colors, "第五部分：落地路径——从0到1", [
        {'type': 'highlight', 'text': '四步走战略'},
        {'type': 'bullet', 'text': '第一步：试点（1-3个月）——选择2-3个高价值场景，内部验证'},
        {'type': 'bullet', 'text': '第二步：扩展（3-6个月）——扩展到10大核心场景，建立数字员工矩阵'},
        {'type': 'bullet', 'text': '第三步：赋能（6-12个月）——数字员工能力对外输出，服务上下游企业'},
        {'type': 'bullet', 'text': '第四步：生态（12-24个月）——构建管网行业AI生态，成为示范企业'},
    ])
    
    # 21. 试点计划
    add_content_slide(prs, colors, "第一步：试点——具体行动计划", [
        {'type': 'highlight', 'text': '试点场景选择标准'},
        {'type': 'table', 'rows': [
            ['标准', '说明', '管网集团应用'],
            ['高价值', '能显著降本增效', '管道巡检、故障预警'],
            ['低风险', '失败影响可控', '报告生成、培训教育'],
            ['易验证', '3个月内可出效果', '报告生成效率提升'],
            ['可复制', '能推广到其他场景', '巡检→维护→安全'],
        ], 'cols': 3},
        {'type': 'highlight', 'text': '试点团队组建'},
        {'type': 'bullet', 'text': '业务专家：管道巡检、故障预警经验'},
        {'type': 'bullet', 'text': '技术专家：AI模型训练、智能体开发'},
        {'type': 'bullet', 'text': '数据专家：管网数据治理、标注'},
        {'type': 'bullet', 'text': '项目经理：跨部门协调、进度管理'},
    ])
    
    # 22. ROI计算
    add_content_slide(prs, colors, "ROI计算——数字员工的经济账", [
        {'type': 'highlight', 'text': '以管道巡检智能体为例'},
        {'type': 'table', 'rows': [
            ['项目', '传统模式', '智能体模式', '差异'],
            ['人力成本', '10人×10万/年=100万', '2人×10万/年=20万', '↓80万'],
            ['巡检效率', '100km/天', '300km/天', '↑200%'],
            ['漏检率', '5%', '0.5%', '↓90%'],
            ['Token成本', '0', '5万/年', '+5万'],
            ['年净节省', '-', '-', '≈75万/年'],
        ], 'cols': 4},
        {'type': 'highlight', 'text': '投资回收期：< 1个月'},
    ])
    
    # 23. Token成本收益
    add_content_slide(prs, colors, "Token计费模式下的成本收益", [
        {'type': 'highlight', 'text': '管网集团数字员工Token消耗预估'},
        {'type': 'table', 'rows': [
            ['智能体', '日均Token消耗', '月均成本', '月均节省', 'ROI'],
            ['管道巡检', '1000万', '5000元', '6.5万', '13x'],
            ['故障预警', '500万', '2500元', '8万', '32x'],
            ['泵站调度', '200万', '1000元', '1.25万', '12.5x'],
            ['报告生成', '100万', '500元', '1万', '20x'],
            ['合计', '1810万', '9000元', '16.75万', '18.6x'],
        ], 'cols': 5},
        {'type': 'highlight', 'text': 'Token成本极低，但节省显著。数字员工的ROI远超传统IT投资。'},
    ])
    
    # 24. 战略启示
    add_content_slide(prs, colors, "战略启示", [
        {'type': 'highlight', 'text': '对国家管网集团甘肃分公司'},
        {'type': 'table', 'rows': [
            ['维度', '建议'],
            ['政策对接', '主动对接国发〔2025〕11号，争取"AI+能源"示范项目'],
            ['场景开放', '开放管网巡检、调度、安全场景，争取考核高分'],
            ['国资考核', '利用"战略价值>ROI"政策，大胆投资AI'],
            ['资金申请', '申请绿电改造、REITs、绿色债券、"两重两新"资金'],
            ['示范引领', '争取成为管网行业AI示范企业，全国推广'],
        ], 'cols': 2},
    ])
    
    # 25. 核心结论
    add_conclusion_slide(prs, colors)
    
    # 26. 行动清单
    add_content_slide(prs, colors, "管网集团行动清单", [
        {'type': 'highlight', 'text': '立即行动（1个月内）'},
        {'type': 'bullet', 'text': '成立AI智能体专项工作组'},
        {'type': 'bullet', 'text': '选择2-3个试点场景（巡检+预警+报告）'},
        {'type': 'bullet', 'text': '对接甘肃省发改委，了解考核要求'},
        {'type': 'bullet', 'text': '启动绿电改造可行性研究'},
        {'type': 'highlight', 'text': '短期行动（3个月内）'},
        {'type': 'bullet', 'text': '完成试点场景验证'},
        {'type': 'bullet', 'text': '建立数字员工矩阵（10大场景）'},
        {'type': 'bullet', 'text': '申请"AI+能源"示范项目'},
        {'type': 'bullet', 'text': '启动REITs/绿色债券申报'},
        {'type': 'highlight', 'text': '中期行动（12个月内）'},
        {'type': 'bullet', 'text': '数字员工能力对外输出'},
        {'type': 'bullet', 'text': '构建管网行业AI生态'},
        {'type': 'bullet', 'text': '争取成为全国示范企业'},
        {'type': 'bullet', 'text': 'Token计费模式商业化'},
    ])
    
    # 27. Q&A
    add_qa_slide(prs, colors)
    
    # 保存
    output_path = os.path.join(os.path.dirname(__file__), "浙大授课-管网集团AI智能体培训.pptx")
    prs.save(output_path)
    print(f"✅ PPTX已生成：{output_path}")
    return output_path

if __name__ == '__main__':
    main()
